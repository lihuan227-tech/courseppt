#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
仰望星空 · Day 5: 我 的 太空 梦想 展示  My Space Dream Showcase
Co-created class booklet: 我们 的 太空 梦想
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random, os

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

NIGHT  = RGBColor(0x0D,0x1B,0x3E)
COSMIC = RGBColor(0x6A,0x1B,0x9A)
STAR   = RGBColor(0xF5,0xC2,0x42)
GOLD   = RGBColor(0xFF,0xB7,0x00)
EARTH  = RGBColor(0x1E,0x88,0xE5)
RED    = RGBColor(0xC8,0x25,0x3E)
MARS   = RGBColor(0xD8,0x43,0x15)
PINK   = RGBColor(0xEC,0x40,0x7A)
NEBULA = RGBColor(0x7B,0x1F,0xA2)
ALIEN  = RGBColor(0x66,0xBB,0x6A)
SKY    = RGBColor(0x42,0xA5,0xF5)
MOON_C = RGBColor(0xB0,0xBE,0xC5)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=IMGBG; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=STAR,a=PP_ALIGN.CENTER)
    for x,y in [(0.8,4.7),(1.8,4.5),(7.8,4.5),(8.6,4.7)]:
        d=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(x),Inches(y),Inches(0.35),Inches(0.35))
        d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
    return s
def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mx=x+size/2; my=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mx-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(my-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

n=0

# Cover
s=ns(); bg(s,NIGHT)
random.seed(31)
for _ in range(50):
    x=random.uniform(0.3,9.7); y=random.uniform(0.3,5.3); sz=random.uniform(0.07,0.18)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
tb(s,0.5,0.4,9,0.6,"🌌 仰望星空  Looking Up at the Stars",sz=28,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.0,9,0.45,"Day 5 · 我 的 太空 梦想 🎉  My Space Dream",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# Big star + planets
sh=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(4.0),Inches(1.85),Inches(2.0),Inches(2.0))
sh.fill.solid(); sh.fill.fore_color.rgb=STAR; sh.line.fill.background()
tb(s,4.0,2.55,2.0,0.6,"🌟",sz=70,a=PP_ALIGN.CENTER)
tb(s,2.0,3.0,1.5,1.0,"🚀",sz=55,a=PP_ALIGN.CENTER)
tb(s,6.5,3.0,1.5,1.0,"🛸",sz=55,a=PP_ALIGN.CENTER)
tb(s,0.5,4.20,9,0.40,"📕 班级 共创 绘本: 我们 的 太空 梦想",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,4.60,9,0.25,"Co-created class book: Our Space Dreams",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.5,4.95,9,0.30,"🎤 Final Showcase — 你 准备好 上台 了 吗?",sz=12,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 5 开场:\n• 「最后 一天! 我们 这周 学了 好多!」\n• 「今天 — 上午 整合 + 共创 绘本, 下午 完成 + 大 Showcase!」")

# Session 1 Divider
s=div("Session 1  上午 11:00–11:45","🌟 故事课 · 我们 的 太空 梦想  ·  45 min",COSMIC,"📕"); n+=1; pn(s,n)

# Learning Goals
s=ns(); bg(s,CREAM); hb(s,"🎯 今天的学习目标  Today's Learning Goals",NIGHT)
tb(s,0.4,0.85,9.2,0.30,"上完这节课, 你会……",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
goals=[
    ("1","📚","通过 共创 绘本 — 整合 一周 学到 的!",NEBULA),
    ("2","💭","表达 你 的 未来 梦想 + 创意 想法",PINK),
    ("3","🗣️","用 简单 中文 介绍 你 的 太空 梦想",STAR),
    ("4","💪","培养 表达 + 自信 + 展示 能力",GOLD),
]
for i,(num,em,text,cl) in enumerate(goals):
    y=1.30+i*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.18),Inches(0.50),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.22,0.50,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.18,0.60,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.90,y+0.20,7.6,0.55,text,sz=14,b=True,c=DARK)
n+=1; pn(s,n)

# Week Recap
s=ns(); bg(s,CREAM); hb(s,"🔁 这 周 我们 学了 什么?  Week Recap",NEBULA)
tb(s,0.4,0.85,9.2,0.32,"4 天 太空 之 旅 — 一起 回顾!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
days=[
    ("D1","☀️","太阳系","Solar System",COSMIC),
    ("D2","⭐","星座 神话","Constellations",PINK),
    ("D3","🚀","太空 探索","Exploration",MARS),
    ("D4","👽","外星人","Aliens",ALIEN),
]
for i,(tag,em,cn,en,cl) in enumerate(days):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.50),Inches(2.22),Inches(3.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=cl; sh.line.fill.background()
    tb(s,x+0.05,1.65,2.12,0.45,tag,sz=18,b=True,c=STAR,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.10,2.12,1.20,em,sz=80,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.40,2.12,0.42,cn,sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.85,2.12,0.30,en,sz=10,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,4.75,9.2,0.30,"🌟 这 一周 — 你 学到 了 好多!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,5.10,9.2,0.24,"This week — you learned SO MUCH!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"4 分钟 — 回顾:\n• 每天 让 一 个 学生 说 一句 「我 记得 ___」")

# Dream prompt
s=ns(); bg(s,NIGHT); hb(s,"💭 你 的 太空 梦想 是 什么?  Your Space Dream?",STAR)
tb(s,0.4,0.85,9.2,0.30,"想 大 一 点! 没有 不 可能!",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Dream BIG! Nothing is impossible!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
dreams=[
    ("👨‍🚀","我 想 当 宇航员!","I want to be an astronaut!",MOON_C),
    ("🚀","我 想 去 火星!","I want to go to Mars!",MARS),
    ("👽","我 想 见 外星人!","I want to meet aliens!",ALIEN),
    ("🏠","我 想 在 月球 建 房子!","I want to build a moon house!",STAR),
    ("📡","我 想 发现 新 星球!","I want to find new planets!",COSMIC),
    ("✏️","我 想 写 太空 故事!","I want to write space stories!",PINK),
]
for i,(em,cn,en,cl) in enumerate(dreams):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.55+row*1.80
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.65))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.10,2.85,0.75,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.85,2.85,0.40,cn,sz=13,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.22,2.85,0.30,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.20),Inches(9.2),Inches(0.35))
prompt.fill.solid(); prompt.fill.fore_color.rgb=GOLD; prompt.line.fill.background()
tb(s,0.55,5.24,9.0,0.28,"💬 「我 的 太空 梦想 是 ___」",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5 分钟:\n• 念 6 个 例子 — 启发 学生\n• 强调: 想 大 一 点! 想 奇怪 一 点!\n• 给 学生 1-2 分钟 自己 想")

# Co-created book overview
s=ns(); bg(s,CREAM); hb(s,"📕 我们 的 太空 梦想 · 班级 共创 绘本",NEBULA)
tb(s,0.4,0.85,9.2,0.32,"每 个 人 写 / 画 一 页 — 合起来 = 我们 的 班级 大 书!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Each student makes 1 page — together = our class book!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 4 page structure
pages=[
    ("1️⃣","我 的 名字 + 画 我 自己","Cover with name + self portrait",STAR),
    ("2️⃣","我 的 太空 梦想 (画 + 写)","My dream (drawing + writing)",NEBULA),
    ("3️⃣","为什么 这 是 我 的 梦想","Why this is my dream",PINK),
    ("4️⃣","10 年 后 — 我 会 ___","In 10 years, I will ___",GOLD),
]
for i,(num,cn,en,cl) in enumerate(pages):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.65))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.15),Inches(y+0.18),Inches(0.85),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+0.15,y+0.25,0.85,0.40,num,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+1.15,y+0.22,3.30,0.42,cn,sz=13,b=True,c=cl)
    tb(s,x+1.15,y+0.66,3.30,0.26,en,sz=9,c=GRAY)
    ib(s,x+1.15,y+0.95,3.30,0.55,"🖍️ 画 这里")
tb(s,0.4,5.10,9.2,0.30,"📕 完成 → 老师 装订 → 这 是 你们 的 班级 礼物!",sz=12,b=True,c=NEBULA,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3 分钟 — 介绍 共创 绘本:\n• 老师 准备: 4 页 模板 (每人 1 套)\n• 完成 后 装订 — 这 是 班级 永久 礼物\n• 鼓励: 「这 是 我们 一起 做 的 大 书!」")

# Session 1 wrap — quick share
s=ns(); bg(s,CREAM); hb(s,"🎤 跟 同桌 分享  Share with Partner",PINK)
tb(s,0.4,0.85,9.2,0.32,"先 跟 同桌 说 — 等下 上台 时 就 不 紧张!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Practice with your partner first — easier on stage!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("1️⃣","你 的 太空 梦想 是 什么?",PINK),
    ("2️⃣","你 为什么 喜欢 这个 梦想?",NEBULA),
    ("3️⃣","你 长大 想 学 什么 让 梦想 成真?",STAR),
]
for i,(num,q,cl) in enumerate(qs):
    y=1.55+i*1.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.00))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,0.55,y+0.30,0.80,0.40,num,sz=24,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.30,8.0,0.50,q,sz=15,b=True,c=cl)
n+=1; pn(s,n)
notes(s,"5 分钟 — Turn & Talk:\n• 2 人 一组\n• 每人 3 个 问题, 各 说 1 分钟\n• 老师 走动 — 听 灵感")

# Session 2 Divider
s=div("Session 2  下午 2:00–2:45","📚 词汇课 · 我会认 + 我会写  ·  45 min",EARTH,"📖"); n+=1; pn(s,n)

# Review
s=ns(); bg(s,CREAM); hb(s,"🔁 早上 学了 什么?",EARTH)
items=[
    ("🌌","一周 回顾","Week recap","太阳系 → 星座 → 探索 → 外星人",NEBULA),
    ("💭","我 的 梦想","My dream","你 想 当 / 想 做 什么?",PINK),
    ("📕","共创 绘本","Class book","每 人 一 页 — 合起来!",STAR),
    ("🤝","同桌 分享","Partner share","练习 — 准备 上台!",GOLD),
]
tb(s,0.4,0.85,9.2,0.32,"想 一 想 — 早上 我们 学了 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,cn,en,d,cl) in enumerate(items):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.40+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.20,1.00,1.20,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.18,3.25,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.20,y+0.58,3.25,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.92,3.25,0.70,d,sz=11,b=True,c=DARK)
n+=1; pn(s,n)

# 我会认 — 5 words
s=ns(); bg(s,CREAM); hb(s,"📖 我会认  I Can Read",STAR)
tb(s,0.4,0.85,9.2,0.32,"5 个 重要 词 — 一起 读!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
words=[
    ("💭","梦想","mèng xiǎng","dream",PINK),
    ("🔮","未来","wèi lái","future",NEBULA),
    ("🔭","探索","tàn suǒ","explore",STAR),
    ("💡","发明","fā míng","invent",GOLD),
    ("🎤","展示","zhǎn shì","showcase",MARS),
]
for i,(em,cn,py,en,cl) in enumerate(words):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.78),Inches(3.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,1.70,0.90,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,1.70,0.55,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.25,1.70,0.35,py,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.62,1.70,0.30,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.05,1.60,0.85,"跟读\n3 遍",sz=11,b=True,c=cl,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

# 我会写 — 梦想 + 未来
def write_slide(emoji,word_cn,word_en,chars,color):
    s=ns(); bg(s,CREAM); hb(s,f"✏️ 我会写 · {word_cn}  I Can Write · {word_en}",color)
    tb(s,0.4,0.85,9.2,0.36,f"{emoji} 一起来写「{word_cn}」!",sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.25,9.2,0.26,f"Practice writing {word_cn} ({word_en})",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tianzi(s,0.55,1.65,2.20,chars[0][0],color,pinyin=chars[0][1],char_sz=120)
    tianzi(s,2.95,1.65,2.20,chars[1][0],color,pinyin=chars[1][1],char_sz=120)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.45,1.72,4.10,0.40,"✏️ 怎么写  How to Write",sz=13,b=True,c=WHITE)
    for i,(ch,py,hint_cn,hint_en) in enumerate(chars):
        y=2.30+i*0.95
        tb(s,5.45,y,4.10,0.35,f"📐「{ch}」 — {py}",sz=13,b=True,c=DARK)
        tb(s,5.45,y+0.32,4.10,0.30,hint_cn,sz=10,b=True,c=color)
        tb(s,5.45,y+0.58,4.10,0.26,hint_en,sz=9,c=GRAY)
    pf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
    pf.fill.solid(); pf.fill.fore_color.rgb=WARM; pf.line.color.rgb=color; pf.line.width=Pt(2)
    tb(s,0.55,4.72,9.0,0.32,f"📝 在 田字格 里 写 3 遍",sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.55,5.08,9.0,0.32,f"💬 「我 的 ___ 是 ___」",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    return s

s=write_slide("💭","梦想","Dream",[
    ("梦","mèng","11 笔 — 「林」+「夕」","Forest + evening (dreams come at night)"),
    ("想","xiǎng","13 笔 — 「相」+「心」","With your heart, you think"),
],PINK); n+=1; pn(s,n)

s=write_slide("🔮","未来","Future",[
    ("未","wèi","5 笔 — 像「木」 上面 加 短横","Like 木 with a short top stroke"),
    ("来","lái","7 笔 — 「未」 + 多 一 横","Like 未 with an extra line"),
],NEBULA); n+=1; pn(s,n)

# Session 3 Divider
s=div("Session 3  下午 3:00–4:30","🎤 项目 + Final Showcase  ·  90 min",GOLD,"🌟"); n+=1; pn(s,n)

# Projects overview
s=ns(); bg(s,CREAM); hb(s,"🛠️ 3 件 事  3 Things to Do",GOLD)
tb(s,0.4,0.85,9.2,0.32,"完成 + 准备 — 然后 大 Showcase!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
projects=[
    ("📕","完成 共创 绘本","Finish class book","每 人 4 页 — 装订 起来",NEBULA),
    ("👨‍🚀","太空 头盔 / 城市","Helmet / City","选 一个 你 想 做!",STAR),
    ("🎤","Final Showcase","Final Show","上台 分享 你 的 梦想!",MARS),
]
for i,(em,cn,en,d,cl) in enumerate(projects):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(3.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.20,em,sz=72,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.85,0.42,cn,sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.38,2.85,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.85,d,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

# Helmet / Space City details
s=ns(); bg(s,CREAM); hb(s,"👨‍🚀🏙️ 太空 头盔 / 太空 城市",STAR)
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(4.10))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=MOON_C; left.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=MOON_C; head.line.fill.background()
tb(s,0.55,0.97,4.30,0.40,"👨‍🚀 太空 头盔 制作",sz=14,b=True,c=DARK)
items=[
    "🌟 你 要 上台 当 宇航员!",
    "1️⃣ 拿 一个 大 纸袋 / 圆 塑料 杯",
    "2️⃣ 画 / 剪 出 大 圆 窗户",
    "3️⃣ 加 天线 / 灯 / 银 色 装饰",
    "4️⃣ 用 锡纸 包 一下 (闪 闪 亮!)",
    "5️⃣ 戴 上 — 上 台 时 戴!",
]
for i,line in enumerate(items):
    tb(s,0.55,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(4.10))
right.fill.solid(); right.fill.fore_color.rgb=WHITE; right.line.color.rgb=NEBULA; right.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,5.20,0.97,4.30,0.40,"🏙️ 太空 城市 设计",sz=14,b=True,c=WHITE)
parts=[
    "🌟 你 的 未来 太空 城市!",
    "1️⃣ 在 大 纸 上 画 整个 城市",
    "2️⃣ 必备: 房子 + 学校 + 游乐场",
    "3️⃣ 加: 太阳能 板 / 温室 / 火车",
    "4️⃣ 旁边 写: 「我 的 城市 叫 ___」",
    "5️⃣ 可以 4 人 一组 一起 做!",
]
for i,line in enumerate(parts):
    tb(s,5.20,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"25 分钟:\n• 头盔: 适合 喜欢 表演 的\n• 城市: 适合 喜欢 设计 的 (可 团队)\n• 材料 提前 准备 好")

# Final Showcase
s=ns(); bg(s,NIGHT); hb(s,"🎤 Final Showcase!  上台 分享!",GOLD)
tb(s,0.4,0.85,9.2,0.30,"轮到 你 上台 — 戴 头盔, 拿 你 的 绘本!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Your turn on stage — wear your helmet, hold your book!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
# Stage area
stage=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(9.2),Inches(2.10))
stage.fill.solid(); stage.fill.fore_color.rgb=WARM; stage.line.color.rgb=STAR; stage.line.width=Pt(3)
tb(s,0.55,1.65,9.0,0.40,"💬 你 的 自我 介绍 (3 句):",sz=13,b=True,c=NEBULA,a=PP_ALIGN.CENTER)
tb(s,0.55,2.10,9.0,0.40,"1️⃣ 「大家 好! 我 是 ___」",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,2.55,9.0,0.40,"2️⃣ 「我 的 太空 梦想 是 ___」",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,3.05,9.0,0.40,"3️⃣ 「因为 ___」",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Bottom — celebration
ce=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.55))
ce.fill.solid(); ce.fill.fore_color.rgb=GOLD; ce.line.color.rgb=STAR; ce.line.width=Pt(3)
tb(s,0.55,3.95,9.0,0.45,"🌟 同学 们 — 给 大家 大 大 的 掌声!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.50,9.0,0.35,"👏 Big applause for everyone!",sz=12,b=True,c=NIGHT,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.30,"📷 拍 合照 + 颁发 「小小 宇航员」 证书!",sz=12,b=True,c=NIGHT,a=PP_ALIGN.CENTER)
tb(s,0.55,5.25,9.0,0.20,"Group photo + 'Little Astronaut' certificate!",sz=8,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"20-25 分钟 — Final Showcase:\n• 每 个 学生 轮流 上台 (30-60 秒/人)\n• 老师 + 家长 鼓掌 (如果 邀请 家长)\n• 最后 — 颁发 证书 + 合照\n• 老师 准备 证书: 「我 是 小小 宇航员! 我 的 梦想 是 ___」")

# Goodbye!
s=ns(); bg(s,NIGHT)
random.seed(99)
for _ in range(60):
    x=random.uniform(0.3,9.7); y=random.uniform(0.3,5.3); sz=random.uniform(0.08,0.20)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
tb(s,0.5,0.6,9,0.8,"🌟 谢谢 大家!  Thank You!",sz=44,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.6,9,0.6,"仰望星空 — 完!",sz=28,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.3,9,0.5,"Looking Up at the Stars — Completed!",sz=18,c=LGRAY,a=PP_ALIGN.CENTER)
# Big star
sh=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(4.2),Inches(3.0),Inches(1.6),Inches(1.6))
sh.fill.solid(); sh.fill.fore_color.rgb=STAR; sh.line.fill.background()
tb(s,4.2,3.30,1.6,1.0,"🚀",sz=70,a=PP_ALIGN.CENTER)
tb(s,0.5,4.85,9,0.40,"💫 继续 仰望 星空 — 你 就 是 未来 的 宇航员!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,5.25,9,0.25,"Keep looking up — YOU are the future astronauts!",sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"结束:\n• 大 合照\n• 发 证书 + 共创 绘本 (装订 好)\n• 「这周 你们 太 棒 了! 继续 仰望 星空!」")

out=os.path.join(os.path.dirname(__file__),"day5_dreams.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
