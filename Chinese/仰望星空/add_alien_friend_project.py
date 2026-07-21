#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Add "Design Your Alien Friend!" PBL project slides to Session 3 of day4_aliens.pptx.
Inserts 5 new slides between current Project 2+3 slide and the Share/Close slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os, copy

SRC = os.path.join(os.path.dirname(__file__), "day4_aliens.pptx")
OUT = SRC  # overwrite in place

prs = Presentation(SRC)
W, H = prs.slide_width, prs.slide_height
LAYOUT = prs.slide_layouts[6]  # OBJECT_WITH_CAPTION_TEXT — matches existing project slides

# Colors (mirrors create_day4_aliens.py)
NIGHT  = RGBColor(0x0D,0x1B,0x3E)
COSMIC = RGBColor(0x6A,0x1B,0x9A)
STAR   = RGBColor(0xF5,0xC2,0x42)
EARTH  = RGBColor(0x1E,0x88,0xE5)
MARS   = RGBColor(0xD8,0x43,0x15)
NEBULA = RGBColor(0x7B,0x1F,0xA2)
ALIEN  = RGBColor(0x66,0xBB,0x6A)
SKY    = RGBColor(0x42,0xA5,0xF5)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
PINK   = RGBColor(0xEC,0x40,0x7A)
ORANGE = RGBColor(0xFB,0x8C,0x00)
TEAL   = RGBColor(0x00,0x89,0x7B)

def ns():
    return prs.slides.add_slide(LAYOUT)

def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt
    r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf

def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)

def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)

def pn(s,n):
    tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)

def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line

def panel(s,l,t,w,h,color,fill=WHITE,lw=2.5):
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    p.fill.solid(); p.fill.fore_color.rgb=fill
    p.line.color.rgb=color; p.line.width=Pt(lw)
    return p

def panel_head(s,l,t,w,color,txt,text_color=WHITE,sz=14):
    h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(0.50))
    h.fill.solid(); h.fill.fore_color.rgb=color; h.line.fill.background()
    tb(s,l+0.15,t+0.07,w-0.3,0.40,txt,sz=sz,b=True,c=text_color)

def chip(s,l,t,w,h,emoji,label,color,sz_em=24,sz_label=11):
    """Rounded chip with emoji + label (e.g., planet choice)."""
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=color; sh.line.width=Pt(2)
    tb(s,l,t+0.05,w,h*0.55,emoji,sz=sz_em,a=PP_ALIGN.CENTER)
    tb(s,l,t+h*0.62,w,h*0.32,label,sz=sz_label,b=True,c=color,a=PP_ALIGN.CENTER)

# ============================================================
# Slide A · Project intro: "Design Your Alien Friend!"
# ============================================================
s=ns(); bg(s,CREAM)
hb(s,"👽 升级 项目 · Design Your Alien Friend!  设计 你 的 外星朋友",ALIEN)

# Hero mission card (left)
panel(s,0.40,0.95,5.40,4.05,ALIEN)
panel_head(s,0.40,0.95,5.40,ALIEN,"🚀 Your Mission  你 的 任务")
tb(s,0.55,1.60,5.10,0.40,"Zorp 有 一个 问题:",sz=13,b=True,c=DARK)
tb(s,0.55,2.05,5.10,0.50,"「宇宙 里 还有 别的 外星人 吗?」",sz=16,b=True,c=COSMIC)
tb(s,0.55,2.55,5.10,0.35,"\"Are there other aliens in space?\"",sz=11,c=GRAY)
tb(s,0.55,3.05,5.10,0.40,"现在 — 轮 到 你 想象!",sz=14,b=True,c=MARS)
tb(s,0.55,3.45,5.10,0.35,"Now it's YOUR turn to imagine!",sz=11,c=GRAY)
tb(s,0.55,3.95,5.10,0.50,"🌟 创造 你 自己 的 外星朋友!",sz=15,b=True,c=ALIEN)
tb(s,0.55,4.42,5.10,0.40,"Create your own alien friend.",sz=11,c=GRAY)

# 4-step preview (right)
panel(s,5.95,0.95,3.65,4.05,COSMIC)
panel_head(s,5.95,0.95,3.65,COSMIC,"🗺️ 4 步 完成  4 Steps")
steps=[
    ("1️⃣","选 星球","Choose a Planet",MARS),
    ("2️⃣","设计 外星人","Design Your Alien",ALIEN),
    ("3️⃣","像 科学家 思考","Think Like a Scientist",EARTH),
    ("4️⃣","动手 做!","Build!",ORANGE),
]
for i,(num,cn,en,cl) in enumerate(steps):
    y=1.60+i*0.78
    tb(s,6.05,y,0.45,0.40,num,sz=18,b=True,c=cl)
    tb(s,6.55,y,3.00,0.35,cn,sz=13,b=True,c=DARK)
    tb(s,6.55,y+0.36,3.00,0.28,en,sz=10,c=GRAY)
pn(s,24)
notes(s,"项目 升级 版 · 90 分钟 总:\n• 这 是 一个 PBL (Project-Based Learning) 升级\n• 学生 通过 4 个 步骤 创造 一个 完整 的 外星朋友\n• 强调: 因为 星球 不同, 外星人 也 不同 — 引入 科学 思维")

# ============================================================
# Slide B · Step 1 (Choose Planet) + Step 2 (Design Alien)
# ============================================================
s=ns(); bg(s,CREAM)
hb(s,"🪐 Step 1 + 2 · 选 星球 + 设计 外星人",MARS)

# Left: Step 1 - Choose a Planet
panel(s,0.40,0.95,4.55,4.05,MARS)
panel_head(s,0.40,0.95,4.55,MARS,"1️⃣ Step 1 · 选 星球  Choose a Planet",sz=13)
tb(s,0.55,1.55,4.30,0.35,"你 的 外星人 住 在 哪里?",sz=12,b=True,c=DARK)
tb(s,0.55,1.88,4.30,0.30,"Where does your alien live?",sz=10,c=GRAY)

planets=[
    ("☀️","Sun","太阳",ORANGE),
    ("🌙","Moon","月球",GRAY),
    ("☄️","Mercury","水星",NEBULA),
    ("🔴","Mars","火星",MARS),
    ("🪐","New Planet","新 星球",ALIEN),
]
# 5 chips in a row at bottom of left panel
chip_w=0.78; gap=0.10
total_w=5*chip_w+4*gap
start_x=0.40+(4.55-total_w)/2
for i,(em,en,cn,cl) in enumerate(planets):
    x=start_x+i*(chip_w+gap)
    chip(s,x,2.40,chip_w,1.05,em,en,cl,sz_em=22,sz_label=9)
    tb(s,x,3.50,chip_w,0.28,cn,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)

tb(s,0.55,4.05,4.30,0.35,"💡 选 一个 — 或 自己 创造 新 星球!",sz=11,b=True,c=MARS,a=PP_ALIGN.CENTER)
tb(s,0.55,4.42,4.30,0.30,"Pick one — or invent your own!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)

# Right: Step 2 - Design Your Alien
panel(s,5.05,0.95,4.55,4.05,ALIEN)
panel_head(s,5.05,0.95,4.55,ALIEN,"2️⃣ Step 2 · 设计 外星人  Design Alien",sz=13)
tb(s,5.20,1.55,4.30,0.32,"想 一 想:  Think about…",sz=11,b=True,c=DARK)

design_qs=[
    ("✅","What is its name?","名字 是 什么?"),
    ("✅","How many eyes?","几 只 眼睛?"),
    ("✅","How many arms?","几 只 手 / 触角?"),
    ("✅","Can it fly / jump / glow?","会 飞? 跳? 发光?"),
]
for i,(mk,en,cn) in enumerate(design_qs):
    y=1.95+i*0.68
    tb(s,5.20,y,0.30,0.32,mk,sz=14,b=True,c=ALIEN)
    tb(s,5.55,y,3.95,0.32,en,sz=11,b=True,c=DARK)
    tb(s,5.55,y+0.32,3.95,0.28,cn,sz=10,c=GRAY)

pn(s,25)
notes(s,"Step 1 · 5 分钟: 选 一个 星球\nStep 2 · 15 分钟: 想象 外星人 的 样子\n• 老师 可以 让 学生 在 纸 上 先 画 草图\n• 鼓励 学生 给 外星人 起 一个 有 创意 的 名字")

# ============================================================
# Slide C · Step 3 (Think Like Scientist) + Step 4 (Build)
# ============================================================
s=ns(); bg(s,CREAM)
hb(s,"🔬 Step 3 + 4 · 像 科学家 思考 + 动手 做!",EARTH)

# Left: Step 3 - Think Like a Scientist
panel(s,0.40,0.95,4.55,4.05,EARTH)
panel_head(s,0.40,0.95,4.55,EARTH,"3️⃣ Step 3 · Think Like a Scientist")
tb(s,0.55,1.55,4.30,0.35,"星球 不一样 — 外星人 也 不一样!",sz=12,b=True,c=COSMIC)
tb(s,0.55,1.88,4.30,0.28,"Different planets = different aliens",sz=10,c=GRAY)

tb(s,0.55,2.25,4.30,0.32,"问 自己:  Ask yourself…",sz=11,b=True,c=DARK)
science_qs=[
    ("🔥","Is it hot there?","那里 热 吗?",MARS),
    ("❄️","Is it cold there?","那里 冷 吗?",SKY),
    ("💨","Is there air?","有 空气 吗?",ALIEN),
    ("👕","Special clothes needed?","要 穿 特别 衣服 吗?",NEBULA),
]
for i,(em,en,cn,cl) in enumerate(science_qs):
    y=2.65+i*0.55
    tb(s,0.55,y,0.40,0.40,em,sz=14)
    tb(s,1.00,y,3.85,0.30,en,sz=10,b=True,c=cl)
    tb(s,1.00,y+0.27,3.85,0.26,cn,sz=9,c=DARK)

# Right: Step 4 - Build!
panel(s,5.05,0.95,4.55,4.05,ORANGE)
panel_head(s,5.05,0.95,4.55,ORANGE,"4️⃣ Step 4 · Build!  动手 做!")
tb(s,5.20,1.55,4.30,0.32,"Craft Materials  材料:",sz=11,b=True,c=DARK)

materials=[
    ("📄","paper","彩纸"),
    ("🥤","cups","纸 杯"),
    ("🎀","pipe cleaners","毛根"),
    ("👁️","googly eyes","活动 眼睛"),
    ("✨","foil","铝箔"),
    ("📦","cardboard","纸 板"),
    ("🟣","pom poms","毛 球"),
    ("🥤","straws","吸管"),
]
# 4 rows × 2 cols
for i,(em,en,cn) in enumerate(materials):
    row=i//2; col=i%2
    x=5.20+col*2.15
    y=1.95+row*0.62
    tb(s,x,y,0.38,0.30,em,sz=13)
    tb(s,x+0.38,y,1.70,0.28,en,sz=9,b=True,c=DARK)
    tb(s,x+0.38,y+0.25,1.70,0.26,cn,sz=9,c=GRAY)

pn(s,26)
notes(s,"Step 3 · 10 分钟: 引导 学生 思考 星球 的 环境 怎样 影响 外星人 的 样子\nStep 4 · 30 分钟: 动手 制作\n• 准备 好 所有 材料 — 让 学生 自由 取用\n• 老师 可以 示范 1-2 个 例子, 但 强调 没有 对错")

# ============================================================
# Slide D · Sentence Frames + Teacher Questions
# ============================================================
s=ns(); bg(s,CREAM)
hb(s,"🎤 分享 时间 · Show Your Alien!",COSMIC)

# Left panel: K-2 frames
panel(s,0.40,0.95,3.00,3.40,SKY)
panel_head(s,0.40,0.95,3.00,SKY,"📘 K-2  低 年级",sz=12)
tb(s,0.55,1.55,2.75,0.30,"Say these:",sz=10,b=True,c=DARK)
k2_lines=[
    "「This is my alien.」",
    "「Its name is ___.」",
    "「It lives on ___.」",
]
for i,line in enumerate(k2_lines):
    tb(s,0.55,1.90+i*0.48,2.75,0.42,line,sz=11,b=True,c=SKY)

# Middle panel: Grade 3-5 frames
panel(s,3.55,0.95,3.00,3.40,NEBULA)
panel_head(s,3.55,0.95,3.00,NEBULA,"📗 Grade 3-5  高 年级",sz=12)
tb(s,3.70,1.55,2.75,0.30,"Add reasons:",sz=10,b=True,c=DARK)
g35_lines=[
    "「This is my alien.」",
    "「Its name is ___.」",
    "「It lives on ___ because ___.」",
    "「It needs ___ because ___.」",
    "「It can ___.」",
]
for i,line in enumerate(g35_lines):
    tb(s,3.70,1.85+i*0.42,2.75,0.40,line,sz=10,b=True,c=NEBULA)

# Right panel: Teacher Q's
panel(s,6.70,0.95,2.90,3.40,ALIEN)
panel_head(s,6.70,0.95,2.90,ALIEN,"❓ 老师 提问  Teacher Q's",sz=12)
teacher_qs=[
    "Why does it live there?",
    "What does it eat?",
    "Can it breathe?",
    "Is it friendly?",
    "Can it come to Earth?",
]
for i,q in enumerate(teacher_qs):
    tb(s,6.85,1.55+i*0.46,2.65,0.40,"• "+q,sz=10,b=True,c=DARK)

# Bottom tip
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(4.50),Inches(9.20),Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb=COSMIC; tip.line.fill.background()
tb(s,0.55,4.58,9.00,0.40,"💬 鼓励 学生 — 多 说 中文 + 英文! Encourage code-switching!",
   sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)

pn(s,27)
notes(s,"分享 · 15 分钟:\n• 每个 学生 1-2 分钟\n• 低 年级 用 简单 句型, 高 年级 加 因果\n• 老师 用 提问 引导 — 让 全班 一起 想!")

# ============================================================
# Slide E · Upgrade — Alien + Home (PBL Extension)
# ============================================================
s=ns(); bg(s,CREAM)
hb(s,"🏘️ Alien + Home  ·  PBL 加 强 版",ORANGE)

# Top intro card
intro=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(0.95),Inches(9.20),Inches(1.00))
intro.fill.solid(); intro.fill.fore_color.rgb=WARM
intro.line.color.rgb=ORANGE; intro.line.width=Pt(2.5)
tb(s,0.55,1.02,9.00,0.35,"💡 想 更 进 一步 (PBL)?  Want to go deeper?",sz=13,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
tb(s,0.55,1.40,9.00,0.30,"不要 只 做 外星人 — 做 「外星人 + 他 的 家」!",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,1.68,9.00,0.26,"Don't just make the alien — make Alien + Home! Different planets need different houses.",
   sz=10,c=GRAY,a=PP_ALIGN.CENTER)

# 3 example house cards
examples=[
    ("🔴","火星 Mars","保暖 house","Warm house\n(很 冷! Cold!)",MARS),
    ("🌙","月球 Moon","Oxygen dome","氧气 罩\n(没 空气! No air!)",SKY),
    ("☀️","太阳 附近","Heat shield house","隔热 房 子\n(太 热! Too hot!)",ORANGE),
]
card_w=2.95; gap=0.18
total=3*card_w+2*gap
start=(10-total)/2
for i,(em,planet,en,cn,cl) in enumerate(examples):
    x=start+i*(card_w+gap)
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.20),Inches(card_w),Inches(2.50))
    p.fill.solid(); p.fill.fore_color.rgb=WHITE
    p.line.color.rgb=cl; p.line.width=Pt(2.5)
    # Planet header strip
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.20),Inches(card_w),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=cl; head.line.fill.background()
    tb(s,x,2.27,card_w,0.40,f"{em}  {planet}",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # Arrow + house type
    tb(s,x,2.85,card_w,0.40,"→",sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x,3.25,card_w,0.40,en,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # Reason (2 lines)
    tb(s,x+0.10,3.75,card_w-0.20,0.85,cn,sz=10,c=GRAY,a=PP_ALIGN.CENTER)

# Bottom punch line
punch=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(4.85),Inches(9.20),Inches(0.55))
punch.fill.solid(); punch.fill.fore_color.rgb=ALIEN; punch.line.fill.background()
tb(s,0.55,4.93,9.00,0.40,"🌟 这样 会 非常 棒!  This makes the project AMAZING!",
   sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)

pn(s,28)
notes(s,"升级 PBL 选项:\n• 适合 高 年级 或 时间 充足 的 班级\n• 让 学生 思考: 星球 环境 → 房子 设计\n• 这是 真正 的 工程师 思维!")

# ============================================================
# Move new slides (currently at end) into position BEFORE the
# original Share/Close slide. They were appended at indices
# 53, 54, 55, 56, 57. We want them at 52..56 (so they sit
# between Project 2+3 and Share/Close).
# Also: bump page number on the now-last slide (Share/Close)
# from 24 → 29 to stay consistent.
# ============================================================
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
total = len(ids)
new_slide_ids = ids[-5:]
share_id = ids[-6]  # was originally last (Share/Close), now sits before the 5 new ones

# Remove the new slide IDs and reinsert them before the Share/Close slide
for sid in new_slide_ids:
    sldIdLst.remove(sid)
# Find current index of share_id and insert before it
new_ids = list(sldIdLst)
share_pos = new_ids.index(share_id)
for offset, sid in enumerate(new_slide_ids):
    sldIdLst.insert(share_pos + offset, sid)

# Update page number on Share/Close slide from 24 -> 29
share_slide = prs.slides[-1]  # after reorder, Share/Close is last
for shape in share_slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip() == "24":
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == "24":
                    run.text = "29"

prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides)} slides total)")
