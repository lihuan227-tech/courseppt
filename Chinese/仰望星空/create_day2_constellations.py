#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
仰望星空 · Day 2: 星空、星座与神话传说  Constellations & Myths
Picture book anchor: 牛郎织女  Cowherd and Weaver Girl
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random, os

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# Palette — Mythic Night
NIGHT  = RGBColor(0x0D,0x1B,0x3E)
COSMIC = RGBColor(0x6A,0x1B,0x9A)
STAR   = RGBColor(0xF5,0xC2,0x42)
GOLD   = RGBColor(0xFF,0xB7,0x00)
EARTH  = RGBColor(0x1E,0x88,0xE5)
RED    = RGBColor(0xC8,0x25,0x3E)
PINK   = RGBColor(0xEC,0x40,0x7A)
NEBULA = RGBColor(0x7B,0x1F,0xA2)
SKY    = RGBColor(0x42,0xA5,0xF5)
SILVER = RGBColor(0xB0,0xBE,0xC5)
GREEN  = RGBColor(0x2E,0x7D,0x32)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=IMGBG; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=STAR,a=PP_ALIGN.CENTER)
    for x,y in [(0.8,4.7),(1.8,4.5),(7.8,4.5),(8.6,4.7),(2.0,1.0),(8.0,1.0)]:
        d=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(x),Inches(y),Inches(0.35),Inches(0.35))
        d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
    return s
def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mx=x+size/2; my=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mx-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(my-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

n=0

# 1. Cover
s=ns(); bg(s,NIGHT)
random.seed(11)
for _ in range(50):
    x=random.uniform(0.3,9.7); y=random.uniform(0.3,5.3); sz=random.uniform(0.06,0.16)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
tb(s,0.5,0.4,9,0.6,"🌌 仰望星空  Looking Up at the Stars",sz=28,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.0,9,0.45,"Day 2 · 星空、星座 与 神话  Constellations & Myths",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# IMAGE PLACEHOLDER — replace with real cover image (e.g., 牛郎织女 / 北斗七星)
ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.5),Inches(1.75),Inches(5.0),Inches(2.85))
ph.fill.solid(); ph.fill.fore_color.rgb=IMGBG; ph.line.color.rgb=STAR; ph.line.width=Pt(3)
tb(s,2.5,2.55,5.0,0.7,"🖼️",sz=60,a=PP_ALIGN.CENTER)
tb(s,2.5,3.30,5.0,0.36,"图片 位置  Image Placeholder",sz=14,b=True,c=NIGHT,a=PP_ALIGN.CENTER)
tb(s,2.5,3.70,5.0,0.28,"老师 后期 放 真实 图片",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,2.5,4.00,5.0,0.28,"Teacher: insert real image later",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.5,4.75,9,0.40,"📖 绘本: 牛郎织女  Cowherd and Weaver Girl",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,5.15,9,0.25,"Picture book — a Chinese love legend in the stars",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 2 开场:\n• 「昨天 我们 学了 太阳系 — 今天 我们 看 星星 和 星座!」\n• 古人 用 故事 解释 星星 — 我们 来 听 一个 中国 故事")

# ============================================================
# 2. Session 1 Divider — 50-min K-5 lesson, 7 parts
# ============================================================
s=div("Session 1  上午 11:00–11:50","⭐ 星座 与 传说 · Constellations & Legends · 50 min",COSMIC,"💫"); n+=1; pn(s,n)

# 3. Learning Goals
s=ns(); bg(s,CREAM); hb(s,"🎯 今天的学习目标  Today's Learning Goals",NIGHT)
tb(s,0.4,0.85,9.2,0.30,"上完这节课, 你会……  By the end, you'll be able to…",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
goals=[
    ("1","⭐","知道 什么 是 「星座」 — 古人 把 星星 连 起来 想象 的 图案",STAR),
    ("2","🔭","知道 古人 为什么 看 星星 (找方向 / 记季节 / 讲故事)",NEBULA),
    ("3","🌟","认识 4 个 经典 星座 + 听 2 个 星座 故事",PINK),
    ("4","✨","创造 属于 你 自己 的 星座!",COSMIC),
]
for i,(num,em,text,cl) in enumerate(goals):
    y=1.30+i*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.18),Inches(0.50),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.22,0.50,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.18,0.60,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.90,y+0.20,7.6,0.55,text,sz=14,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"1 分钟 — 学习目标 (老师 快速 念 一遍, 不细讲)")

# ============================================================
# PART 1 · WARM-UP (5-8 min) — open questions + guess game
# ============================================================
# 4. Warm-up Q1: 你晚上看过星星吗?
s=ns(); bg(s,NIGHT); hb(s,"🌙 Part 1 · 暖身  Warm-up",STAR)
tb(s,0.4,0.85,9.2,0.30,"先 一起 聊 一聊 — 没有 错 答案!",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Let's chat — there are no wrong answers!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
qs=[
    ("🌟","你 晚上 看过 星星 吗? 是 在 哪里 看到 的?","Have you seen stars at night? Where were you?",STAR),
    ("✨","星星 看 起来 像 什么?","What do the stars look like to you?",GOLD),
    ("🔗","如果 把 这些 星星 连 起来, 会 变成 什么?","If you connect them — what shape do you make?",PINK),
]
for i,(em,q,en,cl) in enumerate(qs):
    y=1.55+i*1.20
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,0.55,y+0.30,0.80,0.55,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,1.50,y+0.18,7.95,0.50,q,sz=15,b=True,c=cl)
    tb(s,1.50,y+0.65,7.95,0.30,en,sz=10,c=GRAY)
n+=1; pn(s,n)
notes(s,"⏰ 3-4 分钟 — 暖身:\n• 老师 一题 一题 问 — 收 2-3 个 答案\n• 不 评判 — 接住 每 一个 想法 (写 在 白板)\n• 鼓励 K-1 也 举手 — 「星星 像 我 妈妈 的 耳环!」 都 OK")

# 5. Warm-up — guess game (real photo placeholder of random star points)
s=ns(); bg(s,CREAM); hb(s,"🔍 猜 一 猜 — 你 看到 什么?  Guess Game!",STAR)
tb(s,0.4,0.85,9.2,0.30,"看 这些 星点 — 你 觉得 它们 像 什么?",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Look at these dots — what shape do they make?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# LEFT: real-photo placeholder
ib(s,0.4,1.55,4.55,3.40,"🖼️ 真实 星点 图片 / Random star-points image\n图片 位置 · Image placeholder")
# Right panel with prompt cards
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(3.40))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=STAR; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=STAR; head.line.fill.background()
tb(s,5.25,1.62,4.30,0.40,"💬 想想看  Brainstorm",sz=13,b=True,c=NIGHT)
ideas=[
    "· 像 一只 小狗 吗?",
    "· 像 一个 房子?",
    "· 像 一只 蝴蝶?",
    "· 像 一个 小朋友?",
    "· 像 一棵 树?",
    "· 你 看到 什么? __________",
]
for i,line in enumerate(ideas):
    tb(s,5.30,2.20+i*0.42,4.30,0.40,line,sz=12,b=True,c=DARK)
bottom=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.05),Inches(9.2),Inches(0.45))
bottom.fill.solid(); bottom.fill.fore_color.rgb=COSMIC; bottom.line.color.rgb=STAR; bottom.line.width=Pt(2)
tb(s,0.55,5.10,9.0,0.34,"🙋 不同 的 人 会 看到 不同 东西 — 都 对!",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 3-4 分钟 — 猜 图 互动:\n• 老师 在 左边 放 一张 「随机 星点」 图 (网上 找 — 7-9 颗 散点)\n• 「你 看到 什么?」 — 让 3-4 个 学生 猜\n• 关键: 「同 一组 星 — 不同 人 看到 不一样!」\n• 引出 Part 2: 「这 就 是 星座!」")

# ============================================================
# PART 2 · WHAT IS A CONSTELLATION (8 min)
# ============================================================
# 6. 星座 = ?
s=ns(); bg(s,CREAM); hb(s,"⭐ Part 2 · 什么 是 星座?  What is a Constellation?",COSMIC)
tb(s,0.4,0.85,9.2,0.32,"星座 不是 真的 动物 或 人……",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"A constellation isn't really an animal or a person…",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Big definition card
core=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.60),Inches(9.0),Inches(2.20))
core.fill.solid(); core.fill.fore_color.rgb=COSMIC; core.line.color.rgb=STAR; core.line.width=Pt(3)
tb(s,0.55,1.75,8.9,0.50,"⭐ 星座 = ⭐",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.30,8.9,0.55,"古人 把 星星 连 起来,",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,2.85,8.9,0.55,"想象 出来 的 图案!",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,3.40,8.9,0.30,"Stars connected into a picture in your imagination",sz=11,c=STAR,a=PP_ALIGN.CENTER)
# Bottom: 3 mini illustrations
mini=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.95),Inches(9.0),Inches(1.55))
mini.fill.solid(); mini.fill.fore_color.rgb=WARM; mini.line.color.rgb=COSMIC; mini.line.width=Pt(2)
tb(s,0.6,4.05,8.8,0.32,"💡 像 什么? 看 你 的 想象力!",sz=13,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.7,4.45,2.7,0.35,"⭐⭐⭐",sz=20,a=PP_ALIGN.CENTER)
tb(s,0.7,4.78,2.7,0.30,"3 颗星 → 三角形?",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,3.6,4.45,2.7,0.35,"⭐⭐⭐⭐",sz=20,a=PP_ALIGN.CENTER)
tb(s,3.6,4.78,2.7,0.30,"4 颗星 → 房子? 风筝?",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,6.5,4.45,2.7,0.35,"⭐⭐⭐⭐⭐⭐⭐",sz=18,a=PP_ALIGN.CENTER)
tb(s,6.5,4.78,2.7,0.30,"7 颗星 → 大勺子? 龙?",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 4 分钟 — 解释 「星座」:\n• 念 大字句 2 遍, 全班 跟读\n• 强调: 「不 是 真的 动物 — 是 人 的 想象!」\n• 用 下面 3 组 星点 演示")

# 7. Same dots → different shapes (interactive — each card has real image space)
s=ns(); bg(s,CREAM); hb(s,"🤔 同 一 组 星 — 你 看到 什么?  Same Stars · Different Shapes",NEBULA)
tb(s,0.4,0.85,9.2,0.32,"为什么 不同 的 人 会 看到 不同 的 图案?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Why do different people see different shapes?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Three cards — each has an image placeholder + label
patterns=[
    ("中国 古人","Chinese","北斗七星 → 大勺子","🖼️ 大勺子 图",STAR),
    ("西方 古人","Western","北斗七星 → 大熊 尾巴","🖼️ 大熊座 图",NEBULA),
    ("你 自己?","You?","你 觉得 像 ___?","🖼️ 你 来 画!",PINK),
]
for i,(cn,en,d,img_lb,cl) in enumerate(patterns):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    # Image placeholder area inside card
    ph=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+0.15),Inches(1.70),Inches(2.65),Inches(1.30))
    ph.fill.solid(); ph.fill.fore_color.rgb=IMGBG; ph.line.color.rgb=cl; ph.line.width=Pt(1.5)
    tb(s,x+0.15,2.10,2.65,0.40,img_lb,sz=11,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,2.45,2.65,0.30,"图片 位置",sz=8,c=LGRAY,a=PP_ALIGN.CENTER)
    # Title + caption
    tb(s,x+0.05,3.10,2.85,0.36,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,2.85,0.26,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.55,d,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.55),Inches(9.2),Inches(0.95))
prompt.fill.solid(); prompt.fill.fore_color.rgb=NEBULA; prompt.line.color.rgb=STAR; prompt.line.width=Pt(2.5)
tb(s,0.55,4.62,9.0,0.30,"💬 「我 觉得 这 像 ___」 — 你 的 想象力 是 你 的!",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.30,"I think it looks like ___ — your imagination is yours!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.55,5.25,9.0,0.20,"💡 没有 标准 答案 — 都 是 对 的",sz=9,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 4 分钟 — 同 星 不同 图:\n• 老师 解释: 中国 看 大勺子, 西方 看 大熊\n• 问 学生: 「你 觉得 像 什么?」\n• 重点: 想象力 没 错 答案!")

# ============================================================
# PART 3 · WHY ANCIENT PEOPLE WATCHED STARS (5-7 min)
# ============================================================
# 8. 古时候 没有……
s=ns(); bg(s,CREAM); hb(s,"🕰️ Part 3 · 为什么 古人 看 星星?  Why Ancient People Watched Stars",GOLD)
tb(s,0.4,0.85,9.2,0.32,"想想 看 — 古时候 (1000 年 以前) 没有 这些:",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Long ago (1000+ years ago) — they didn't have any of this:",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
no_things=[
    ("📱","手机","Phone",RED),
    ("🗺️","GPS / 地图","GPS / Map",NEBULA),
    ("💡","电灯","Electric light",GOLD),
    ("📺","电视","TV",COSMIC),
]
for i,(em,cn,en,cl) in enumerate(no_things):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.22),Inches(2.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.65,2.12,0.85,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.55,2.12,0.40,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.12,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # Big red X
    tb(s,x+0.05,3.20,2.12,0.40,"❌",sz=22,a=PP_ALIGN.CENTER)
# Bottom: but they had stars!
bot=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.65))
bot.fill.solid(); bot.fill.fore_color.rgb=NIGHT; bot.line.color.rgb=STAR; bot.line.width=Pt(3)
tb(s,0.55,3.95,9.0,0.40,"……但是 他们 有 星星! ⭐",sz=18,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.40,9.0,0.30,"…but they had the stars!",sz=11,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,4.75,9.0,0.32,"用 星星 → 找 方向 · 记 季节 · 讲 故事 · 想象 天空",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,5.10,9.0,0.26,"Use stars to find direction · track seasons · tell stories · imagine the sky",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 3 分钟 — 古时候 没 这些:\n• 一个 一个 念 — 学生 「啊?!」\n• 关键: 「没有 这些 — 他们 怎么 办?」\n• 引导: 抬头 看 星星!")

# 9. Interactive: 没有 GPS 怎么 办?
s=ns(); bg(s,CREAM); hb(s,"🤔 没有 Google Maps, 你 怎么 回家?  No GPS — How Would YOU Get Home?",NEBULA)
tb(s,0.4,0.85,9.2,0.32,"想 一想 — 古人 是 怎么 找 路 的?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"How did ancient people find their way?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
ways=[
    ("🌟","抬头 找 星星","Look up at stars","北极星 一直 在 北边!",STAR),
    ("🌅","看 太阳","Watch the sun","太阳 从 东边 升起",GOLD),
    ("🌳","记 路 上 的 树 / 山","Remember trees / mountains","看到 大 山 = 快 到 家",NEBULA),
    ("🗣️","问 路上 的 人","Ask someone","大人 知道 路!",PINK),
]
for i,(em,cn,en,d,cl) in enumerate(ways):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.60))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,x+0.10,y+0.12,1.00,0.85,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.10,3.25,0.40,cn,sz=14,b=True,c=cl)
    tb(s,x+1.20,y+0.50,3.25,0.28,en,sz=9,c=GRAY)
    tb(s,x+1.20,y+0.85,3.25,0.65,d,sz=11,b=True,c=DARK)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.10),Inches(9.2),Inches(0.40))
prompt.fill.solid(); prompt.fill.fore_color.rgb=GOLD; prompt.line.fill.background()
tb(s,0.55,5.15,9.0,0.30,"💬 你 还 能 想到 什么 办法?",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 3-4 分钟 — 互动:\n• 让 学生 想象 — 「没有 手机, 你 怎么 回家?」\n• 收 3-5 个 答案\n• 引导: 「最 重要 的 一个 是 — 看 星星!」")

# ============================================================
# PART 4 · 4 CLASSIC CONSTELLATIONS (10-12 min)
# ============================================================
# 10. Overview of 4 constellations — each card has its own image placeholder
s=ns(); bg(s,CREAM); hb(s,"🌟 Part 4 · 4 个 经典 星座  4 Classic Constellations",STAR)
tb(s,0.4,0.85,9.2,0.30,"今天 我们 认识 4 个 — 都 在 你 家 的 天空!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"4 constellations — all visible from your backyard!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
overview=[
    ("北斗七星","Big Dipper","像 大 勺子",STAR),
    ("猎户座","Orion","像 猎人",GOLD),
    ("天鹅座","Cygnus","像 飞 天鹅",PINK),
    ("狮子座","Leo","像 狮子",NEBULA),
]
for i,(cn,en,d,cl) in enumerate(overview):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.22),Inches(3.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    # Image placeholder (top of card)
    ph=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+0.12),Inches(1.70),Inches(1.98),Inches(1.85))
    ph.fill.solid(); ph.fill.fore_color.rgb=IMGBG; ph.line.color.rgb=cl; ph.line.width=Pt(1.5)
    tb(s,x+0.12,2.30,1.98,0.40,"🖼️",sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.12,2.78,1.98,0.30,"图片 位置",sz=10,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.12,3.10,1.98,0.24,"Image",sz=8,c=LGRAY,a=PP_ALIGN.CENTER)
    # Caption
    tb(s,x+0.05,3.65,2.12,0.40,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.05,2.12,0.28,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.40,2.12,0.42,d,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 1 分钟 — 4 个 星座 总览:\n• 老师 在 每张 卡 上 放 一张 真实 星座 图\n• 念 4 个 名字 + 4 个 形状\n• 「下面 我们 一个 一个 看!」")

# 11. 北斗七星 detail
s=ns(); bg(s,CREAM); hb(s,"🥄 北斗七星  Big Dipper",STAR)
tb(s,0.4,0.85,9.2,0.30,"7 颗星 — 像 一个 大 勺子!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"7 stars — looks like a big spoon (Big Dipper)!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Image placeholder for picture (left)
ib(s,0.4,1.55,4.5,2.95,"🖼️ 北斗七星 图片 / Image placeholder")
# Right: facts + question
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(2.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=STAR; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=STAR; head.line.fill.background()
tb(s,5.25,1.62,4.30,0.40,"⭐ 关于 北斗七星",sz=13,b=True,c=NIGHT)
items=[
    "🥄 像 一个 大 勺子",
    "🎯 帮 你 找 北极星 (一直 指 北边)",
    "🌌 全年 都 能 看到 它",
    "🐻 西方 人 觉得 它 是 大熊 的 尾巴",
]
for i,line in enumerate(items):
    tb(s,5.25,2.20+i*0.55,4.30,0.45,line,sz=12,b=True,c=DARK)
ask=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
ask.fill.solid(); ask.fill.fore_color.rgb=COSMIC; ask.line.color.rgb=STAR; ask.line.width=Pt(2)
tb(s,0.55,4.72,9.0,0.32,"💬 你 觉得 真的 像 大 勺子 吗?",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,5.08,9.0,0.30,"Does it really look like a spoon to you?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 2-3 分钟 — 北斗七星:\n• 展示 真实 图片 (老师 准备)\n• 「7 颗星 — 数 一数!」\n• 关键 互动: 「真的 像 勺子 吗?」 — 让 学生 自由 说")

# 12. 猎户座 detail
s=ns(); bg(s,CREAM); hb(s,"🏹 猎户座  Orion",GOLD)
tb(s,0.4,0.85,9.2,0.30,"3 颗 腰带 星 — 像 一个 大 猎人!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"3 belt stars — looks like a big hunter!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.4,1.55,4.5,2.95,"🖼️ 猎户座 图片 / Image placeholder")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(2.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=GOLD; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=GOLD; head.line.fill.background()
tb(s,5.25,1.62,4.30,0.40,"⭐ 关于 猎户座",sz=13,b=True,c=NIGHT)
items=[
    "🏹 像 一个 大 猎人",
    "✨ 看 中间 — 3 颗 星 一排 (腰带!)",
    "❄️ 冬天 晚上 最 容易 看到",
    "🌌 全世界 都 看 得到",
]
for i,line in enumerate(items):
    tb(s,5.25,2.20+i*0.55,4.30,0.45,line,sz=12,b=True,c=DARK)
ask=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
ask.fill.solid(); ask.fill.fore_color.rgb=COSMIC; ask.line.color.rgb=GOLD; ask.line.width=Pt(2)
tb(s,0.55,4.72,9.0,0.32,"💬 你 觉得 这个 猎人 在 追 什么?",sz=14,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.55,5.08,9.0,0.30,"What do you think the hunter is chasing?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 2-3 分钟 — 猎户座:\n• 强调 中间 的 3 颗 腰带 星 — 最容易 找!\n• 互动: 让 学生 想象 — 猎人 在 追 什么?")

# 13. 天鹅座 detail
s=ns(); bg(s,CREAM); hb(s,"🦢 天鹅座  Cygnus",PINK)
tb(s,0.4,0.85,9.2,0.30,"像 一只 飞 起来 的 大 天鹅!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Looks like a swan flying through the sky!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.4,1.55,4.5,2.95,"🖼️ 天鹅座 图片 / Image placeholder")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(2.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=PINK; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=PINK; head.line.fill.background()
tb(s,5.25,1.62,4.30,0.40,"⭐ 关于 天鹅座",sz=13,b=True,c=WHITE)
items=[
    "🦢 像 一只 大 鸟 张开 翅膀",
    "✨ 也 叫 「北十字」 (像 一个 十字)",
    "☀️ 夏天 晚上 看 最 清楚",
    "💫 就 在 银河 里 飞!",
]
for i,line in enumerate(items):
    tb(s,5.25,2.20+i*0.55,4.30,0.45,line,sz=12,b=True,c=DARK)
ask=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
ask.fill.solid(); ask.fill.fore_color.rgb=COSMIC; ask.line.color.rgb=PINK; ask.line.width=Pt(2)
tb(s,0.55,4.72,9.0,0.32,"💬 你 觉得 它 飞 去 哪里?",sz=14,b=True,c=PINK,a=PP_ALIGN.CENTER)
tb(s,0.55,5.08,9.0,0.30,"Where do you think the swan is flying to?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 2 分钟 — 天鹅座:\n• 让 学生 张 开 手臂 — 当 天鹅\n• 「就 在 银河 里 飞 — 多 美!」")

# 14. 狮子座 detail
s=ns(); bg(s,CREAM); hb(s,"🦁 狮子座  Leo",NEBULA)
tb(s,0.4,0.85,9.2,0.30,"像 一只 大 狮子 — 春天 出现!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Looks like a big lion — appears in spring!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.4,1.55,4.5,2.95,"🖼️ 狮子座 图片 / Image placeholder")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(2.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=NEBULA; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,5.25,1.62,4.30,0.40,"⭐ 关于 狮子座",sz=13,b=True,c=WHITE)
items=[
    "🦁 像 一只 大 狮子 趴着",
    "👑 有 一个 「狮子 头」 + 长 尾巴",
    "🌸 春天 晚上 看 最 清楚",
    "✨ 它 是 12 星座 之 一",
]
for i,line in enumerate(items):
    tb(s,5.25,2.20+i*0.55,4.30,0.45,line,sz=12,b=True,c=DARK)
ask=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
ask.fill.solid(); ask.fill.fore_color.rgb=COSMIC; ask.line.color.rgb=NEBULA; ask.line.width=Pt(2)
tb(s,0.55,4.72,9.0,0.32,"💬 你 觉得 真的 像 狮子 吗? 学一 个 狮子 叫 看!",sz=13,b=True,c=NEBULA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.08,9.0,0.30,"Does it really look like a lion? Try a lion roar!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 2-3 分钟 — 狮子座:\n• 让 学生 学 狮子 叫 — 大家 笑!\n• 「春天 你 抬头 — 找 找 它」")

# ============================================================
# PART 5 · LEGENDS (10 min)
# ============================================================
# 15. 牛郎织女 (Chinese myth)
s=ns(); bg(s,CREAM); hb(s,"📖 Part 5 · 牛郎织女  Cowherd & Weaver",PINK)
tb(s,0.4,0.85,9.2,0.36,"中国 故事 — 关于 银河 和 七夕",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"A Chinese story about the Milky Way and Qixi Festival",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.5,1.65,4.5,3.30,"📚 牛郎织女 / Picture book")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(3.30))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=PINK; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=PINK; head.line.fill.background()
tb(s,5.40,1.72,4.10,0.40,"📖 故事",sz=13,b=True,c=WHITE)
parts=[
    "👦 牛郎 — 一个 放 牛 的 小伙子",
    "👧 织女 — 天上 王母 的 孙女",
    "💕 他们 相爱, 结 婚, 有 了 孩子",
    "😢 王母 用 银河 把 他们 分开!",
    "🐦 每年 七月初七, 喜鹊 搭 桥",
    "✨ 让 他们 一年 见 一次!",
]
for i,line in enumerate(parts):
    tb(s,5.40,2.25+i*0.42,4.10,0.40,line,sz=11,b=True,c=DARK)
ask=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.05),Inches(9.2),Inches(0.45))
ask.fill.solid(); ask.fill.fore_color.rgb=PINK; ask.line.fill.background()
tb(s,0.55,5.10,9.0,0.34,"💬 为什么 他们 一年 只 能 见 一次? 你 是 喜鹊, 你 会 帮 吗?",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 5 分钟 — 牛郎织女:\n• 老师 戏剧 语气 讲\n• 中间 停 一下 — 「啊! 王母 太 坏 了!」\n• 关键 互动 (2 个 问题):\n  1. 「为什么 一年 只 见 一次?」\n  2. 「你 是 喜鹊 — 你 会 帮 吗?」")

# 16. 银河 + 牛郎星 + 织女星 — real Milky Way photo placeholder
s=ns(); bg(s,CREAM); hb(s,"💫 银河 · 牛郎星 · 织女星  Milky Way",COSMIC)
tb(s,0.4,0.85,9.2,0.30,"夏天 晚上 抬头 — 你 真的 能 看到 这 3 颗星!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"On summer nights — you can really see these 3!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Real Milky Way photo placeholder
ib(s,0.4,1.55,9.2,3.00,"🖼️ 真实 银河 + 牛郎星 + 织女星 照片\nReal photo: Milky Way · Altair · Vega · 图片 位置")
# Caption strip — labels for the 3 elements
cap=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
cap.fill.solid(); cap.fill.fore_color.rgb=WHITE; cap.line.color.rgb=COSMIC; cap.line.width=Pt(2.5)
tb(s,0.55,4.72,3.0,0.28,"💫 银河 Milky Way",sz=12,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,3.55,4.72,3.0,0.28,"⭐ 牛郎星 Altair",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,6.55,4.72,3.0,0.28,"⭐ 织女星 Vega",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,5.05,9.0,0.36,"🌌 7 月 7 日 — 喜鹊 帮 他们 搭 桥!",sz=13,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 2 分钟 — 银河 + 两颗 星:\n• 老师 在 中间 放 一张 真实 银河 + 标 牛郎星 + 织女星 的 照片 (网上 找)\n• 指着 照片: 「左边 是 牛郎, 右边 是 织女, 中间 是 银河」\n• 「今晚 抬头 — 找 银河!」")

# 17. 猎户座 Orion (Western myth)
s=ns(); bg(s,CREAM); hb(s,"📖 西方 故事 · 猎户座 Orion",GOLD)
tb(s,0.4,0.85,9.2,0.36,"古希腊 传说 — 一个 勇敢 的 猎人",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"Ancient Greek legend — a brave hunter",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.5,1.65,4.5,3.30,"📚 Orion / Picture")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(3.30))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=GOLD; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=GOLD; head.line.fill.background()
tb(s,5.40,1.72,4.10,0.40,"📖 故事",sz=13,b=True,c=NIGHT)
parts=[
    "🏹 古人 觉得 天上 有 一个 大 猎人",
    "💪 他 又 高 又 强 — 谁 都 不 怕",
    "🐂 他 带着 弓 + 箭 + 大棒",
    "🦂 但 一只 蝎子 螫 了 他!",
    "🌟 神 把 他 放 到 天上 当 星座",
    "🏃 现在 他 还 在 追 狮子 + 公牛",
]
for i,line in enumerate(parts):
    tb(s,5.40,2.25+i*0.42,4.10,0.40,line,sz=11,b=True,c=DARK)
ask=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.05),Inches(9.2),Inches(0.45))
ask.fill.solid(); ask.fill.fore_color.rgb=GOLD; ask.line.fill.background()
tb(s,0.55,5.10,9.0,0.34,"💬 你 觉得 他 在 追 什么? 你 会 帮 哪 一边?",sz=12,b=True,c=NIGHT,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 3 分钟 — Orion:\n• 简单 版 — 不 讲 太多 神 名\n• 关键: 「天空 有 一个 大 猎人 — 一直 在 追 动物!」\n• 互动: 「你 觉得 他 在 追 什么?」")

# ============================================================
# PART 6 · CREATIVE ACTIVITY (10-12 min) — Create your own constellation
# ============================================================
# 18. 创造你的星座!
s=ns(); bg(s,COSMIC); hb(s,"✨ Part 6 · 创造 你 的 星座!  Create YOUR Constellation!",STAR,t=0.15)
tb(s,0.4,0.85,9.2,0.40,"如果 你 可以 创造 一个 属于 你 自己 的 星座 ……",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"If you could make a constellation just for yourself…",sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
# 3 question cards
qs=[
    ("🤔","你 会 创造 什么?","What would you create?"),
    ("🎨","它 长 什么 样?","What does it look like?"),
    ("📖","背后 有 什么 故事?","What's the story behind it?"),
]
for i,(em,cn,en) in enumerate(qs):
    x=0.4+i*3.10; y=1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(2.60))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.20,2.85,1.10,em,sz=72,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.45,2.85,0.45,cn,sz=15,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.95,2.85,0.32,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3 tier instructions
tier=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.45),Inches(9.2),Inches(1.05))
tier.fill.solid(); tier.fill.fore_color.rgb=WARM; tier.line.color.rgb=STAR; tier.line.width=Pt(2.5)
tb(s,0.55,4.52,9.0,0.30,"🎯 不同 年级 不同 任务  Different grades, different tasks",sz=12,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.55,4.85,3.0,0.30,"K-1 · 画 你 的 星座",sz=11,b=True,c=PINK,a=PP_ALIGN.CENTER)
tb(s,3.7,4.85,2.7,0.30,"2-3 · 画 + 一句话",sz=11,b=True,c=NEBULA,a=PP_ALIGN.CENTER)
tb(s,6.5,4.85,3.0,0.30,"4-5 · 画 + 短 故事",sz=11,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.55,5.18,9.0,0.20,"K-1 draw · 2-3 draw + sentence · 4-5 draw + short story",sz=8,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 1-2 分钟 — 介绍 任务:\n• 念 3 个 大 问题 — 让 学生 想\n• 强调 不同 年级 不同 要求 — 不会 写字 也 OK!\n• 老师 准备: 黑色卡纸 + 白色 贴纸 / 棉签 / 颜料")

# 19. Sentence frames + example
s=ns(); bg(s,CREAM); hb(s,"💬 句型 + 例子  Sentence Frames + Example",PINK)
tb(s,0.4,0.85,9.2,0.32,"用 这些 句子 帮 你 — 不会 写 也 没 关系!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Use these sentence frames — drawing alone is OK too!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Left: sentence frames
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(4.55),Inches(3.45))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=COSMIC; left.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=COSMIC; head.line.fill.background()
tb(s,0.55,1.62,4.30,0.40,"💬 句型  Sentence Frames",sz=13,b=True,c=STAR)
frames=[
    "🌟 我 的 星座 叫 ___",
    "🎨 它 看 起来 像 ___",
    "💖 它 代表 ___",
    "📖 它 的 故事 是 ___",
]
for i,line in enumerate(frames):
    tb(s,0.55,2.20+i*0.65,4.30,0.50,line,sz=14,b=True,c=DARK)
# Right: example
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.55),Inches(4.55),Inches(3.45))
right.fill.solid(); right.fill.fore_color.rgb=WARM; right.line.color.rgb=STAR; right.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.55),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=STAR; head.line.fill.background()
tb(s,5.20,1.62,4.30,0.40,"⭐ 老师 的 例子  Example",sz=13,b=True,c=NIGHT)
ex=[
    "🌟 我 的 星座 叫 「小猫座」",
    "🎨 它 看 起来 像 我 家 的 小 黑 猫",
    "💖 它 代表 「想家」",
    "📖 它 的 故事 是: 一只 小猫",
    "       想 妈妈 想 到 飞 上 天 ……",
]
for i,line in enumerate(ex):
    tb(s,5.20,2.20+i*0.50,4.30,0.42,line,sz=12,b=True,c=DARK)
do=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.10),Inches(9.2),Inches(0.40))
do.fill.solid(); do.fill.fore_color.rgb=PINK; do.line.fill.background()
tb(s,0.55,5.15,9.0,0.30,"🎨 8-10 分钟 — 自己 画 + 写 — 老师 走动 帮忙",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 8-10 分钟 — 创作:\n• 老师 念 自己 的 例子 (小猫座) — 让 学生 笑 + 觉得「我 也 行」\n• 材料: 黑色 卡纸 + 白色 贴纸 / 棉签\n• 老师 走动 — 帮 K-1 写字, 鼓励 4-5 写 短 故事\n• 留 2-3 分钟 给 Part 7 分享")

# ============================================================
# PART 7 · SHARE & CLOSE (5 min)
# ============================================================
# 20. Share circle
s=ns(); bg(s,CREAM); hb(s,"🎤 Part 7 · 分享 你 的 星座!  Share Your Constellation!",STAR)
tb(s,0.4,0.85,9.2,0.32,"轮到 你 上台 — 给 同学 介绍 你 的 星座!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Your turn on stage — share your constellation with the class!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Big share frame
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(9.2),Inches(2.40))
sh.fill.solid(); sh.fill.fore_color.rgb=NIGHT; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,0.55,1.65,9.0,0.40,"💬 用 句型 介绍:",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.10,9.0,0.45,"「我 的 星座 叫 ___」",sz=18,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.55,9.0,0.45,"「它 看 起来 像 ___」",sz=18,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,3.00,9.0,0.45,"「它 的 故事 是 ___」",sz=18,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,3.55,9.0,0.30,"👏 同学 鼓掌 — 没有 错 答案!",sz=12,b=True,c=PINK,a=PP_ALIGN.CENTER)
# Closing
close=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.10),Inches(9.2),Inches(1.40))
close.fill.solid(); close.fill.fore_color.rgb=COSMIC; close.line.color.rgb=STAR; close.line.width=Pt(2.5)
tb(s,0.55,4.20,9.0,0.36,"🌟 老师 总结:",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.60,9.0,0.40,"「古人 看 星星 讲 故事 ——",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,5.05,9.0,0.40,"今天 我们 也 创造 了 自己 的 故事!」",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 3-4 分钟 — 分享:\n• 让 4-5 个 学生 上 台 (轮流, 每人 ~30 秒)\n• 鼓励 高低 年级 都 参与\n• 老师 总结: 「古人 + 我们 都 在 看 星星 + 讲 故事!」")

# 21. Wrap question
s=ns(); bg(s,NIGHT); hb(s,"🌙 今晚 抬头 看 看 天空!  Look Up Tonight!",STAR)
tb(s,0.4,0.85,9.2,0.32,"今晚 回家 — 你 想 找 哪 一颗 星星?",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Tonight when you go home — which star will you look for?",sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
# Star sprinkles
import random as _r2
_r2.seed(99)
for _ in range(40):
    x=_r2.uniform(0.3,9.7); y=_r2.uniform(1.6,4.4); sz=_r2.uniform(0.07,0.18)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
# Centered prompt
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.0),Inches(2.4),Inches(6.0),Inches(1.50))
prompt.fill.solid(); prompt.fill.fore_color.rgb=COSMIC; prompt.line.color.rgb=STAR; prompt.line.width=Pt(3)
tb(s,2.0,2.55,6.0,0.50,"我 想 找 ___ ⭐",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,2.0,3.10,6.0,0.40,"I want to look for ___",sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,2.0,3.55,6.0,0.30,"(北斗 / 织女 / 牛郎 / 我 的 星座……)",sz=11,b=True,c=PINK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.95,9.2,0.40,"💫 下次 见 (Day 3): 🚀 太空 探索 + 火箭 实验!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,5.30,9.2,0.22,"Next time (Day 3): Space exploration + rocket experiments!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⏰ 1 分钟 — 收 尾:\n• 让 学生 说 一句 — 「今晚 我 想 找 ___」\n• 「家庭 作业」: 今晚 抬头 看 天空 (跟 爸 妈 一起)\n• 提示 Day 3 是 火箭 + 太空 探索!")

# 8. Session 2 Divider
s=div("Session 2  下午 2:00–2:45","📚 词汇课 · 我会认 + 我会写  ·  45 min",EARTH,"📖"); n+=1; pn(s,n)

# 9. Review
s=ns(); bg(s,CREAM); hb(s,"🔁 早上 学了 什么?  Morning Review",EARTH)
items=[
    ("👦👧","牛郎织女","Cowherd & Weaver","古老 中国 故事",PINK),
    ("⭐","星座","Constellations","把 星星 连 起来 看 形状",STAR),
    ("🌌","银河","Milky Way","一条 白色 大 河",COSMIC),
    ("💡","星空","Starry sky","= 科学 + 故事 + 想象!",NEBULA),
]
tb(s,0.4,0.85,9.2,0.32,"想 一 想 — 早上 我们 学了 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"What did we explore this morning?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
for i,(em,cn,en,d,cl) in enumerate(items):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.20,1.00,1.20,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.18,3.25,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.20,y+0.58,3.25,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.92,3.25,0.70,d,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"2-3 分钟 — 复习")

# 10-14. 我会认 — 5 vocabulary words (one per slide, matching Day 1 format)
read_words=[
    ("⭐","星星","xīng xing","Stars",STAR,
        "晚上 天上 有 好多 星星。",
        "📷 夜空 / 闪亮 的 星"),
    ("✨","星座","xīng zuò","Constellation",GOLD,
        "北斗七星 是 一 个 星座。",
        "📷 大熊座 / 连线 图"),
    ("🌌","银河","yín hé","Milky Way",COSMIC,
        "银河 像 一条 白色 的 河。",
        "📷 银河 / 夏夜 星空"),
    ("👻","神话","shén huà","Myth",PINK,
        "牛郎织女 是 一 个 中国 神话。",
        "📷 牛郎织女 / 喜鹊 桥"),
    ("📖","故事","gù shi","Story",NEBULA,
        "古人 用 故事 解释 星空。",
        "📷 绘本 / 老人 讲 故事"),
]
for em,cn,py,en,c,sent,img_label in read_words:
    s=ns(); bg(s,CREAM); hb(s,f"👀 我会认 · {cn}  I Can Read",c)
    # Left: big character card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.6))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.fill.background()
    tb(s,0.5,1.05,4.3,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,0.5,1.75,4.3,1.0,cn,sz=66 if len(cn)==2 else 72,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,f"{py}  ·  {en}",sz=18,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.25,4.3,0.34,"👉 跟我读!  Read after me!",sz=13,b=True,c=c,a=PP_ALIGN.CENTER)
    # Right: image placeholder
    ib_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.3),Inches(1.0),Inches(4.4),Inches(2.6))
    ib_box.fill.solid(); ib_box.fill.fore_color.rgb=IMGBG; ib_box.line.color.rgb=c; ib_box.line.width=Pt(2)
    tb(s,5.3,1.80,4.4,0.6,"🖼️",sz=44,a=PP_ALIGN.CENTER)
    tb(s,5.3,2.50,4.4,0.4,img_label,sz=12,c=LGRAY,a=PP_ALIGN.CENTER)
    tb(s,5.3,2.95,4.4,0.30,"图片 位置 · Image placeholder",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
    # Bottom: example sentence
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.30))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=c; sh2.line.width=Pt(2.5)
    tb(s,0.6,3.95,2.0,0.40,"📌 例句  Example",sz=14,b=True,c=c)
    tb(s,0.6,4.40,8.8,0.55,sent,sz=22,b=True,c=DARK)
    tb(s,0.4,5.25,9.2,0.28,"💬 「我 认识 ___」  · I know the word ___",sz=11,b=True,c=c,a=PP_ALIGN.CENTER)
    n+=1; pn(s,n)
    notes(s,f"3 分钟 — {cn}:\n• 老师 指 字, 全班 齐读 3 遍 (慢 → 快 → 大声)\n• 看 图: 「这 是 ___, 你 见过 吗?」\n• 读 例句, 学生 跟读\n• 抽 1-2 个 学生 用「{cn}」造 一 个 新 句子\n• 写 到 黑板 上 — 让 学生 在 空中 跟着 写 一遍")

# 11-12. 我会写 — 星星 + 星座
def write_slide(emoji,word_cn,word_en,chars,color):
    s=ns(); bg(s,CREAM); hb(s,f"✏️ 我会写 · {word_cn}  I Can Write · {word_en}",color)
    tb(s,0.4,0.85,9.2,0.36,f"{emoji} 一起来写「{word_cn}」!",sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.25,9.2,0.26,f"Practice writing {word_cn} ({word_en})",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    if len(chars)==2:
        tianzi(s,0.55,1.65,2.20,chars[0][0],color,pinyin=chars[0][1],char_sz=120)
        tianzi(s,2.95,1.65,2.20,chars[1][0],color,pinyin=chars[1][1],char_sz=120)
    else:
        tianzi(s,1.30,1.65,2.95,chars[0][0],color,pinyin=chars[0][1],char_sz=160)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.45,1.72,4.10,0.40,"✏️ 怎么写  How to Write",sz=13,b=True,c=WHITE)
    for i,(ch,py,hint_cn,hint_en) in enumerate(chars):
        y=2.30+i*0.95
        tb(s,5.45,y,4.10,0.35,f"📐「{ch}」 — {py}",sz=13,b=True,c=DARK)
        tb(s,5.45,y+0.32,4.10,0.30,hint_cn,sz=10,b=True,c=color)
        tb(s,5.45,y+0.58,4.10,0.26,hint_en,sz=9,c=GRAY)
    pf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
    pf.fill.solid(); pf.fill.fore_color.rgb=WARM; pf.line.color.rgb=color; pf.line.width=Pt(2)
    tb(s,0.55,4.72,9.0,0.32,f"📝 在 田字格 里 写 3 遍",sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.55,5.08,9.0,0.32,f"💬 「我 会 写「{word_cn}」!」",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    return s

s=write_slide("⭐","星星","Stars",[
    ("星","xīng","9 笔 — 上面「日」+ 下面「生」","Sun on top + 生 below"),
    ("星","xīng","写 两遍 一样 — 重复 词","Repeat character"),
],STAR); n+=1; pn(s,n)
notes(s,"5-6 分钟 — 写 星星")

s=write_slide("✨","星座","Constellation",[
    ("星","xīng","9 笔 — 上「日」+ 下「生」","Sun + life"),
    ("座","zuò","10 笔 — 「广」+「人人」+「土」","Roof + two people + earth"),
],GOLD); n+=1; pn(s,n)
notes(s,"5-6 分钟 — 写 星座")

# 13. Session 3 Divider
s=div("Session 3  下午 3:00–4:30","🛠️ 项目课 · 星座 创作  ·  90 min",NEBULA,"⭐"); n+=1; pn(s,n)

# 14. Projects overview
s=ns(); bg(s,CREAM); hb(s,"🛠️ 3 个 项目  3 Projects",NEBULA)
tb(s,0.4,0.85,9.2,0.32,"选 一个 — 自己 来 当 星座 创造者!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
projects=[
    ("🔦","星座 手电筒 投影","Constellation Flashlight","纸杯 + 针孔 → 投影 到 墙上!",STAR),
    ("🧵","棉签 星座 拼图","Q-tip Constellations","棉签 + 黑纸 — 拼 你 的 星座",NEBULA),
    ("✍️","创造 你 的 星座 故事","Make YOUR Constellation","起 个 名字 + 写 个 故事",PINK),
]
for i,(em,cn,en,d,cl) in enumerate(projects):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(3.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.20,em,sz=78,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.85,0.42,cn,sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.38,2.85,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.85,d,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.90,9.2,0.30,"⭐ 完成 后 — 给 你 的 星座 取 名字!",sz=12,b=True,c=NEBULA,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3 分钟 — 介绍 项目:\n• 项目 1 (手电筒): 需要 提前 准备 纸杯 + 针\n• 项目 2 (棉签): 棉签 + 黑色卡纸\n• 项目 3 (写故事): 适合 喜欢 写字 的 学生")

# 15. Project 1 details — Constellation flashlight
s=ns(); bg(s,CREAM); hb(s,"🔦 项目 1 · 星座 手电筒  Constellation Flashlight",STAR)
ib(s,0.4,0.90,4.5,3.95,"🖼️ 手电筒 投影 示例")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(3.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=STAR; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=STAR; head.line.fill.background()
tb(s,5.25,0.97,4.30,0.40,"📝 怎么 做  How to Make",sz=14,b=True,c=DARK)
steps=[
    ("1️⃣","拿 一个 纸杯 (paper cup)"),
    ("2️⃣","在 杯底 用 针 戳 小孔 — 摆 成 星座"),
    ("3️⃣","可以 是 大熊座 / 北斗七星 / 自创"),
    ("4️⃣","把 手电筒 放进 杯里 — 照 墙上!"),
    ("5️⃣","关 灯 — 看 你 的 星座 在 墙上 发光!"),
]
for i,(num,txt) in enumerate(steps):
    y=1.55+i*0.62
    tb(s,5.25,y,0.40,0.40,num,sz=16,b=True,c=STAR)
    tb(s,5.75,y+0.04,3.80,0.35,txt,sz=11,c=DARK)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.95),Inches(9.2),Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb=NIGHT; tip.line.fill.background()
tb(s,0.55,5.02,9.0,0.32,"💡 安全 提示: 老师 帮 戳 孔 — 不要 自己 用 针!",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"项目 1 · 25 分钟:\n• 材料: 纸杯 + 针 (老师 操作) + 手电筒 (or 手机 闪光灯)\n• 提前 准备: 大熊座 / 牛郎织女 / 北斗 的 图案 模板\n• 课堂 关 灯 试 — 学生 会 惊叹!")

# 16. Project 2 + 3 (combined)
s=ns(); bg(s,CREAM); hb(s,"🧵✍️ 项目 2 + 3  Two More Options",NEBULA)
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(4.10))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=NEBULA; left.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,0.55,0.97,4.30,0.40,"🧵 项目 2 · 棉签 星座 拼图",sz=14,b=True,c=WHITE)
items=[
    "🎨 材料: 黑色卡纸 + 棉签 + 白色 颜料 / 贴纸",
    "1️⃣ 想 — 你 要 拼 什么 星座?",
    "2️⃣ 用 白色 圆 贴纸 贴 星星 位置",
    "3️⃣ 用 棉签 蘸 白颜料 — 连 起来",
    "4️⃣ 写 上 星座 的 名字",
    "✨ 给 你 的 队 / 妈妈 / 朋友!",
]
for i,line in enumerate(items):
    tb(s,0.55,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(4.10))
right.fill.solid(); right.fill.fore_color.rgb=WHITE; right.line.color.rgb=PINK; right.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=PINK; head.line.fill.background()
tb(s,5.20,0.97,4.30,0.40,"✍️ 项目 3 · 创造 你 的 星座 故事",sz=14,b=True,c=WHITE)
parts=[
    "🌟 假装 你 是 古人 — 抬头 看 星星",
    "1️⃣ 把 几 颗 星 连 起来 — 看 像 什么?",
    "2️⃣ 给 它 取 一个 名字 (你的 名字 都 行!)",
    "3️⃣ 编 一个 小 故事:",
    "「很久 以前, 有 一个 ___」",
    "4️⃣ 画 + 写 — 老师 帮 K-1!",
]
for i,line in enumerate(parts):
    tb(s,5.20,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"25 分钟:\n• 项目 2 适合 喜欢 手工 的 学生\n• 项目 3 适合 喜欢 创作 的 学生\n• 老师 走动 — 鼓励 每个 想法")

# 17. Share & close
s=ns(); bg(s,CREAM); hb(s,"🎤 分享 + 再见!",COSMIC)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(2.10))
sh.fill.solid(); sh.fill.fore_color.rgb=NIGHT; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,0.55,1.05,9.0,0.40,"💬 句型 — 分享 你 的 星座!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,1.55,9.0,0.50,"「我 的 星座 叫 ___」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.10,9.0,0.50,"「它 看 起来 像 ___」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
preview=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.25),Inches(9.2),Inches(2.05))
preview.fill.solid(); preview.fill.fore_color.rgb=WHITE; preview.line.color.rgb=COSMIC; preview.line.width=Pt(2.5)
tb(s,0.55,3.40,9.0,0.40,"🔮 下次 见 (Day 3):",sz=14,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.55,3.85,9.0,0.40,"🚀 登月 + 火星 计划 + 火箭 实验!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.30,9.0,0.30,"Moon landing + Mars + rocket experiments",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,4.70,9.0,0.30,"👋 今晚 抬头 — 找 找 银河 + 牛郎织女!",sz=11,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5 分钟 — 分享:\n• 3-4 学生 用 句型 分享 自己 的 星座\n• 「家庭 作业」: 今晚 看 银河")

out=os.path.join(os.path.dirname(__file__),"day2_constellations.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
