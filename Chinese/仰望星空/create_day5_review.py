#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
仰望星空 · Day 5: 总 复习 + 最 后 项目 / Final Review + Showcase

Plan (per user):
  Session 1 (上午, 45 min):  复习 Day 1+2 + 分组 比赛
  Session 2 (下午, 45 min):  复习 Day 3+4 + 分组 比赛
  Session 3 (下午, 90 min):  最 后 项目 — 4 个 分层 项目 (individual or group)

Design constraints:
  • 不太费老师 — minimal teacher prep, students self-directed
  • 分层 — projects span K-2 (low writing) to grade 3-5 (full sentences)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# Palette (consistent with days 1-4)
NIGHT  = RGBColor(0x0D,0x1B,0x3E)
COSMIC = RGBColor(0x6A,0x1B,0x9A)
STAR   = RGBColor(0xF5,0xC2,0x42)
GOLD   = RGBColor(0xFF,0xB7,0x00)
EARTH  = RGBColor(0x1E,0x88,0xE5)
MARS   = RGBColor(0xD8,0x43,0x15)
NEBULA = RGBColor(0x7B,0x1F,0xA2)
ALIEN  = RGBColor(0x66,0xBB,0x6A)
SKY    = RGBColor(0x42,0xA5,0xF5)
PINK   = RGBColor(0xEC,0x40,0x7A)
ORANGE = RGBColor(0xFB,0x8C,0x00)
TEAL   = RGBColor(0x00,0x89,0x7B)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])

def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt
    r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf

def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)

def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)

def pn(s,n):
    tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)

def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line

def panel(s,l,t,w,h,color,fill=WHITE,lw=2.5):
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    p.fill.solid(); p.fill.fore_color.rgb=fill
    p.line.color.rgb=color; p.line.width=Pt(lw)
    return p

def panel_head(s,l,t,w,color,txt,text_color=WHITE,sz=14):
    h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(0.50))
    h.fill.solid(); h.fill.fore_color.rgb=color; h.line.fill.background()
    tb(s,l+0.15,t+0.07,w-0.3,0.40,txt,sz=sz,b=True,c=text_color)

def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,0.3,1.50,9.4,1.0,f"{emoji} {title}",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.3,2.70,9.4,0.4,sub,sz=18,b=True,c=STAR,a=PP_ALIGN.CENTER)
    for x,y in [(0.8,4.7),(1.8,4.5),(7.8,4.5),(8.6,4.7)]:
        d=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(x),Inches(y),Inches(0.35),Inches(0.35))
        d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
    return s

def vocab_chip(s,l,t,w,h,cn,py,en,color,sz_cn=16):
    """Small vocab card: char on top, pinyin + English below."""
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=color; sh.line.width=Pt(2)
    tb(s,l,t+0.05,w,h*0.42,cn,sz=sz_cn,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,l,t+h*0.50,w,h*0.22,py,sz=7,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,l,t+h*0.72,w,h*0.25,en,sz=8,b=True,c=DARK,a=PP_ALIGN.CENTER)

n = 0

# ============================================================
# Slide 1 · COVER
# ============================================================
s = ns(); bg(s, NIGHT)
# stars
import random
random.seed(42)
for _ in range(60):
    x = random.uniform(0.2, 9.7); y = random.uniform(0.2, 5.4); sz = random.choice([0.05, 0.08, 0.10])
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
# Title
tb(s,0.5,1.20,9.0,0.6,"🌌 仰望星空  Looking Up at the Stars",sz=24,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.85,9.0,0.5,"Day 5 · 总 复习 + 最 后 项目",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.35,9.0,0.4,"Final Review + Showcase",sz=18,c=LGRAY,a=PP_ALIGN.CENTER)
# Sub-emoji row
tb(s,0.5,3.20,9.0,0.6,"🪐 ⭐ 🚀 👽 🎉",sz=44,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,4.10,9.0,0.4,"5 天 太空 之 旅 — 一起 庆祝!",sz=16,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,4.50,9.0,0.4,"5-day space journey — let's celebrate!",sz=12,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 5 总 时长: 3 课时 · 大约 3 小时\nSession 1 (45 min): D1+D2 复习 + 比赛\nSession 2 (45 min): D3+D4 复习 + 比赛\nSession 3 (90 min): 4 个 分层 项目 · 选 一 个\n所有 项目 都 是 学生 自主 完成 — 老师 只 需要 准备 材料 + 巡视")

# ============================================================
# Slide 2 · LEARNING GOALS
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎯 今天 的 学习 目标  Today's Learning Goals",STAR)
tb(s,0.4,0.85,9.2,0.32,"上完这节课, 你会……  By the end, you'll be able to…",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)

goals = [
    ("1️⃣","回顾 5 天 学过 的 太空 知识","Recall what you learned this week",NEBULA),
    ("2️⃣","和 同学 一起 玩 太空 知识 比赛","Play knowledge games as a team",ORANGE),
    ("3️⃣","完成 一个 「太空 大 项目」, 带 回家!","Finish a final project to take home",ALIEN),
]
for i,(num,cn,en,cl) in enumerate(goals):
    y = 1.45 + i*1.15
    p = panel(s,0.5,y,9.0,0.95,cl,fill=WHITE,lw=2.5)
    tb(s,0.65,y+0.10,0.6,0.50,num,sz=22,b=True,c=cl)
    tb(s,1.30,y+0.10,8.0,0.40,cn,sz=15,b=True,c=DARK)
    tb(s,1.30,y+0.50,8.0,0.32,en,sz=11,c=GRAY)
n+=1; pn(s,n)
notes(s,"目标 是 复习 + 应用 + 创造 — 不 学 新 词\n如果 时间 紧, 项目 时间 优先 — 比赛 可以 简化")

# ============================================================
# Slide 3 · SESSION 1 DIVIDER
# ============================================================
s=div("Session 1","🪐 上午 45 min · 复习 Day 1+2 + 比赛 大 PK",ALIEN,"🌞"); n+=1; pn(s,n)

# ============================================================
# Slide 4 · Day 1 + Day 2 quick review (combined, side-by-side)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔁 快速 回顾 · Day 1 + Day 2",NEBULA)

# LEFT panel — Day 1 太阳系
panel(s,0.40,0.95,4.55,4.05,GOLD)
panel_head(s,0.40,0.95,4.55,GOLD,"🌞 Day 1 · 太阳系  Solar System",sz=13)
# Mnemonic
tb(s,0.55,1.55,4.30,0.32,"🪐 八 大 行星 口诀:",sz=12,b=True,c=DARK)
tb(s,0.55,1.88,4.30,0.40,"水 · 金 · 地 · 火 · 木 · 土 · 天 · 海",sz=14,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
# Vocab chips
tb(s,0.55,2.45,4.30,0.30,"📚 我 会 认:",sz=11,b=True,c=DARK)
d1_vocab = [
    ("太阳","tài yáng","Sun",GOLD),
    ("月亮","yuè liàng","Moon",SKY),
    ("地球","dì qiú","Earth",EARTH),
    ("星球","xīng qiú","Planet",NEBULA),
    ("宇宙","yǔ zhòu","Universe",COSMIC),
]
chip_w=0.84; gap=0.04
total = 5*chip_w + 4*gap
start_x = 0.40 + (4.55 - total)/2
for i,(cn,py,en,cl) in enumerate(d1_vocab):
    vocab_chip(s, start_x+i*(chip_w+gap), 2.80, chip_w, 1.15, cn, py, en, cl, sz_cn=18)
tb(s,0.55,4.10,4.30,0.30,"📖 故事: 神奇 校车 · 迷失 在 太阳系",sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.42,4.30,0.30,"Magic School Bus · Lost in Solar System",sz=8,c=GRAY,a=PP_ALIGN.CENTER)

# RIGHT panel — Day 2 星座
panel(s,5.05,0.95,4.55,4.05,NEBULA)
panel_head(s,5.05,0.95,4.55,NEBULA,"⭐ Day 2 · 星座  Constellations",sz=13)
tb(s,5.20,1.55,4.30,0.32,"⭐ 古人 看 星星 — 连 成 图案!",sz=12,b=True,c=DARK)
tb(s,5.20,1.88,4.30,0.30,"Ancient people connected stars into pictures",sz=9,c=GRAY)
# Key concepts
star_facts = [
    ("🦁","狮子座 · 天鹅座 · 猎犬座","Constellation shapes"),
    ("🧭","北极星 — 引路 的 星","North star — guides travelers"),
    ("🌌","北斗 七星 — 像 一 把 勺子","Big Dipper — looks like a spoon"),
]
for i,(em,cn,en) in enumerate(star_facts):
    y = 2.30 + i*0.45
    tb(s,5.20,y,0.40,0.40,em,sz=14)
    tb(s,5.60,y,3.95,0.28,cn,sz=10,b=True,c=DARK)
    tb(s,5.60,y+0.26,3.95,0.24,en,sz=8,c=GRAY)

tb(s,5.20,3.75,4.30,0.30,"📚 我 会 认:",sz=11,b=True,c=DARK)
d2_vocab = [
    ("星星","xīng xīng","Star",STAR),
    ("星座","xīng zuò","Constellation",NEBULA),
    ("银河","yín hé","Milky Way",EARTH),
]
chip_w2 = 1.25; gap2 = 0.15
total2 = 3*chip_w2 + 2*gap2
start_x2 = 5.05 + (4.55 - total2)/2
for i,(cn,py,en,cl) in enumerate(d2_vocab):
    vocab_chip(s, start_x2+i*(chip_w2+gap2), 4.08, chip_w2, 0.95, cn, py, en, cl, sz_cn=20)

n+=1; pn(s,n)
notes(s,"5-10 分钟 快速 复习:\n• 学生 看 着 卡片 一起 读 词\n• 老师 问 1-2 个 问题 引出 记忆\n• 不 是 教 新 内容 — 只 是 唤醒")

# ============================================================
# Slide 5 · GAME 1 — D1+D2 比赛 大 PK
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🏆 比赛 1 · 太阳系 + 星座 大 PK",ORANGE)

# Game rules box
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(0.95),Inches(9.20),Inches(0.85))
intro.fill.solid(); intro.fill.fore_color.rgb=WARM
intro.line.color.rgb=ORANGE; intro.line.width=Pt(2.5)
tb(s,0.55,1.02,9.00,0.32,"🎮 规则  Rules · 分 成 2-3 队, 每队 起 一个 太空 队名",sz=12,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
tb(s,0.55,1.35,9.00,0.30,"Form teams · pick a space team name · earn ⭐ for each correct answer",sz=10,c=GRAY,a=PP_ALIGN.CENTER)

# 3 rounds in row
rounds = [
    ("🃏","Round 1","词语 分类",
     "Sort the Vocab","老师 念 词, 队员 举 「D1」或「D2」 牌",
     "Teacher says word, team holds up Day 1 or Day 2 card",GOLD),
    ("✔️❌","Round 2","真 假 题",
     "True or False","老师 念 句子, 对 = 拍手, 错 = 跺脚",
     "Statement → clap if true, stomp if false",NEBULA),
    ("🖼️","Round 3","看 图 猜 词",
     "Picture → Word","老师 画 / 比 划, 全队 喊 词",
     "Teacher draws or mimes, team shouts the word",ALIEN),
]
card_w=2.95; gap=0.12
total=3*card_w + 2*gap; start=(10-total)/2
for i,(em,rd,cn,en,desc_cn,desc_en,cl) in enumerate(rounds):
    x=start+i*(card_w+gap)
    panel(s,x,1.95,card_w,3.05,cl,lw=2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(card_w),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=cl; head.line.fill.background()
    tb(s,x,2.02,card_w,0.40,rd,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x,2.55,card_w,0.40,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x,3.10,card_w,0.32,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x,3.42,card_w,0.28,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.85,card_w-0.20,0.70,desc_cn,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.55,card_w-0.20,0.40,desc_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)

tb(s,0.40,5.10,9.20,0.30,"⭐ 每 答 对 1 题 = 1 颗 星 · 最 多 星 的 队 拿 「太空 冠军」!",sz=12,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30-35 分钟:\n• Round 1 (10 min) — 老师 准备 D1/D2 词语 列表; 学生 举 手势 / 颜色 牌\n• Round 2 (10 min) — 真 假 题 准备 列表 见 老师 备 课 笔记\n• Round 3 (15 min) — 老师 抽 词, 让 学生 上 来 表演 或 画\n\n真 假 题 示例:\n• 太阳 是 行星. (假 — 是 恒星)\n• 北极星 一直 在 北方. (真)\n• 月球 上 有 空气. (假)\n• 银河 里 有 很多 星星. (真)")

# ============================================================
# Slide 6 · SESSION 2 DIVIDER
# ============================================================
s=div("Session 2","🚀 下午 45 min · 复习 Day 3+4 + 比赛 大 PK",MARS,"👽"); n+=1; pn(s,n)

# ============================================================
# Slide 7 · Day 3 + Day 4 quick review
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔁 快速 回顾 · Day 3 + Day 4",COSMIC)

# LEFT panel — Day 3 航天
panel(s,0.40,0.95,4.55,4.05,MARS)
panel_head(s,0.40,0.95,4.55,MARS,"🚀 Day 3 · 探秘 航天  Space Exploration",sz=12)

d3_facts = [
    ("👨‍🚀","宇航员 — 严格 训练!","Astronauts train hard"),
    ("🛰️","太空站 = 太空 里 的 「家」","Space station = home in space"),
    ("🍱","压缩袋装 食物 + 检查 设备","Packaged food + experiments"),
]
for i,(em,cn,en) in enumerate(d3_facts):
    y = 1.55 + i*0.50
    tb(s,0.55,y,0.40,0.40,em,sz=14)
    tb(s,0.95,y,3.95,0.30,cn,sz=11,b=True,c=DARK)
    tb(s,0.95,y+0.28,3.95,0.24,en,sz=8,c=GRAY)

tb(s,0.55,3.15,4.30,0.30,"📚 我 会 认:",sz=11,b=True,c=DARK)
d3_vocab = [
    ("火箭","huǒ jiàn","Rocket",MARS),
    ("月球","yuè qiú","Moon",GRAY),
    ("火星","huǒ xīng","Mars",MARS),
    ("太空站","tài kōng zhàn","Station",SKY),
]
chip_w3 = 1.05; gap3 = 0.07
total3 = 4*chip_w3 + 3*gap3
start_x3 = 0.40 + (4.55 - total3)/2
for i,(cn,py,en,cl) in enumerate(d3_vocab):
    vocab_chip(s, start_x3+i*(chip_w3+gap3), 3.48, chip_w3, 1.05, cn, py, en, cl, sz_cn=18)
tb(s,0.55,4.60,4.30,0.32,"📖 故事: Part 3 · 内 行星",sz=9,b=True,c=DARK,a=PP_ALIGN.CENTER)

# RIGHT panel — Day 4 外星人
panel(s,5.05,0.95,4.55,4.05,ALIEN)
panel_head(s,5.05,0.95,4.55,ALIEN,"👽 Day 4 · 外星人?  Do Aliens Exist?",sz=12)

d4_facts = [
    ("✉️","Zorp 火星人 写信 给 地球","Zorp wrote from Mars"),
    ("🔴","火星 很 冷 · 没 空气","Mars is cold, no air"),
    ("💭","星球 不同 → 外星人 不同","Different planets, different aliens"),
]
for i,(em,cn,en) in enumerate(d4_facts):
    y = 1.55 + i*0.50
    tb(s,5.20,y,0.40,0.40,em,sz=14)
    tb(s,5.60,y,3.95,0.30,cn,sz=11,b=True,c=DARK)
    tb(s,5.60,y+0.28,3.95,0.24,en,sz=8,c=GRAY)

tb(s,5.20,3.15,4.30,0.30,"📚 我 会 认:",sz=11,b=True,c=DARK)
d4_vocab = [
    ("外星人","wài xīng rén","Alien",ALIEN),
    ("生命","shēng mìng","Life",NEBULA),
    ("信号","xìn hào","Signal",SKY),
    ("发现","fā xiàn","Discover",GOLD),
    ("猜想","cāi xiǎng","Guess",COSMIC),
]
chip_w4 = 0.84; gap4 = 0.04
total4 = 5*chip_w4 + 4*gap4
start_x4 = 5.05 + (4.55 - total4)/2
for i,(cn,py,en,cl) in enumerate(d4_vocab):
    # 外星人 is 3 chars — use smaller size to fit on one line
    sz_cn = 13 if len(cn) >= 3 else 18
    vocab_chip(s, start_x4+i*(chip_w4+gap4), 3.48, chip_w4, 1.05, cn, py, en, cl, sz_cn=sz_cn)
tb(s,5.20,4.60,4.30,0.32,"📖 故事: 来自 火星 的 一 封 信",sz=9,b=True,c=DARK,a=PP_ALIGN.CENTER)

n+=1; pn(s,n)

# ============================================================
# Slide 8 · GAME 2 — D3+D4 大 比赛
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🏆 比赛 2 · 航天 + 外星人 大 PK",MARS)

intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(0.95),Inches(9.20),Inches(0.85))
intro.fill.solid(); intro.fill.fore_color.rgb=WARM
intro.line.color.rgb=MARS; intro.line.width=Pt(2.5)
tb(s,0.55,1.02,9.00,0.32,"🎮 规则  Rules · 继续 上 午 的 分队 · 每队 起 一个 太空 队名",sz=12,b=True,c=MARS,a=PP_ALIGN.CENTER)
tb(s,0.55,1.35,9.00,0.30,"Same teams · earn ⭐ for each correct answer · winners get a galaxy badge!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)

rounds2 = [
    ("🔗","Round 1","词语 配 对",
     "Match the Word","卡片 上 词 + emoji, 找 一 对",
     "Match the word card to the picture card",SKY),
    ("📖","Round 2","故事 接 龙",
     "Story Chain","老师 起 头, 学生 接 一 句",
     "Teacher starts a sentence, team adds next part",EARTH),
    ("🖼️","Round 3","画 一 画 猜 一 猜",
     "Draw and Guess","一 人 画, 全队 猜 词",
     "One draws · team shouts the word (Pictionary)",PINK),
]
card_w=2.95; gap=0.12
total=3*card_w + 2*gap; start=(10-total)/2
for i,(em,rd,cn,en,desc_cn,desc_en,cl) in enumerate(rounds2):
    x=start+i*(card_w+gap)
    panel(s,x,1.95,card_w,3.05,cl,lw=2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(card_w),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=cl; head.line.fill.background()
    tb(s,x,2.02,card_w,0.40,rd,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x,2.55,card_w,0.40,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x,3.10,card_w,0.32,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x,3.42,card_w,0.28,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.85,card_w-0.20,0.70,desc_cn,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.55,card_w-0.20,0.40,desc_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)

tb(s,0.40,5.10,9.20,0.30,"🏅 累计 上午 + 下午 总 星 数 · 颁 「太空 总 冠军」奖!",sz=12,b=True,c=MARS,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30-35 分钟:\n• Round 1 (10 min): 把 D3+D4 词 + emoji 印 在 卡片 上, 学生 配对\n• Round 2 (10 min): 老师 起 头, 例如 「Zorp 是 一个 ___」 学生 续 写\n• Round 3 (15 min): 像 Pictionary, 一个 学生 上 来 画, 队友 猜\n\n词 库 (画一画 抽词 用):\n外星人 · 生命 · 信号 · 火箭 · 月球 · 火星 · 太空站 · 宇航员 · 发现 · 猜想")

# ============================================================
# Slide 9 · SESSION 3 DIVIDER
# ============================================================
s=div("Session 3","🎨 下午 90 min · 最 后 项目 · 4 个 选 一 个!",COSMIC,"🚀"); n+=1; pn(s,n)

# ============================================================
# Slide 10 · 4 PROJECTS OVERVIEW
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🛠️ 4 个 项目 · 选 一 个  4 Projects · Pick One",COSMIC)

tb(s,0.4,0.85,9.2,0.32,"💡 每 个 项目 都 有 「难度」 标 — 选 适合 你 的!",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.28,"Each project shows its level — pick what fits you best!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)

projects = [
    ("📖","项目 1","太空 翻翻书","Space Flip Book",
     "K-2","个人 / Individual",
     "画 + 贴 — 把 一 周 学 的 词 做 成 小 书",
     "Draw + sticker — make a mini word book",
     ALIEN),
    ("🃏","项目 2","知识 竞赛 卡","Trivia Card Game",
     "All","小组 / Group",
     "做 10 张 问答 卡, 和 别 的 队 玩",
     "Make 10 Q&A cards · trade and play",
     ORANGE),
    ("📔","项目 3","我 的 太空 日记","My Space Journal",
     "3-5","个人 / Individual",
     "4 页 — 一 天 一 页, 写 + 画",
     "4 pages · one per day · write + draw",
     EARTH),
    ("🎭","项目 4","太空 小 剧场","Space Mini Theater",
     "All","小组 / Group",
     "排 一个 1-2 分钟 太空 小 短剧",
     "Plan + perform a 1-2 min space skit",
     PINK),
]

card_w = 2.20; gap = 0.10
total = 4*card_w + 3*gap; start = (10 - total)/2
for i,(em,pn_label,cn_title,en_title,level,style,desc_cn,desc_en,cl) in enumerate(projects):
    x = start + i*(card_w+gap)
    panel(s,x,1.55,card_w,3.45,cl,lw=2.5)
    # head
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(card_w),Inches(0.45))
    head.fill.solid(); head.fill.fore_color.rgb=cl; head.line.fill.background()
    tb(s,x,1.60,card_w,0.35,pn_label,sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # emoji
    tb(s,x,2.05,card_w,0.55,em,sz=36,a=PP_ALIGN.CENTER)
    # title
    tb(s,x+0.05,2.62,card_w-0.10,0.30,cn_title,sz=12,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.92,card_w-0.10,0.25,en_title,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    # level + style badges
    badge1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.10),Inches(3.25),Inches(0.85),Inches(0.28))
    badge1.fill.solid(); badge1.fill.fore_color.rgb=cl; badge1.line.fill.background()
    tb(s,x+0.10,3.28,0.85,0.22,level,sz=8,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    badge2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+1.00),Inches(3.25),Inches(card_w-1.10),Inches(0.28))
    badge2.fill.solid(); badge2.fill.fore_color.rgb=WHITE
    badge2.line.color.rgb=cl; badge2.line.width=Pt(1)
    tb(s,x+1.00,3.28,card_w-1.10,0.22,style,sz=7,b=True,c=cl,a=PP_ALIGN.CENTER)
    # description
    tb(s,x+0.10,3.65,card_w-0.20,0.65,desc_cn,sz=9,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.30,card_w-0.20,0.55,desc_en,sz=7,c=GRAY,a=PP_ALIGN.CENTER)

n+=1; pn(s,n)
notes(s,"老师 解释 3-5 分钟:\n• 每 个 项目 都 让 学生 自由 选\n• 难度 标 是 建议, 不 是 限制\n• K-2 项目 也 适合 高 年级 — 关键 是 让 他们 享受\n• 老师 准备 各 项目 材料 摆 在 桌 上, 学生 取 用\n\n材料 准备 清单:\n• Project 1: A4 纸 (每人 1 张) + 贴 纸 + 彩笔\n• Project 2: 索引 卡 (每组 15-20 张) + 笔\n• Project 3: 4 页 装订 小 本子 (或 折 纸) + 笔\n• Project 4: 简单 头饰 / 道具 (纸 + 彩笔 自 做)")

# ============================================================
# Slide 11 · Project 1 · 太空 翻翻书 (K-2)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📖 项目 1 · 太空 翻翻书  Space Flip Book",ALIEN)

# Level badge top-right
lvl=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.30),Inches(0.95),Inches(2.30),Inches(0.42))
lvl.fill.solid(); lvl.fill.fore_color.rgb=ALIEN; lvl.line.fill.background()
tb(s,7.30,0.99,2.30,0.34,"🌱 K-2 · 个人 Individual",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)

# LEFT — steps
panel(s,0.40,1.50,4.55,3.50,ALIEN)
panel_head(s,0.40,1.50,4.55,ALIEN,"📝 怎么 做  Steps",sz=12)
steps=[
    "1️⃣ 拿 一 张 A4 纸 — 折 成 8 格",
    "2️⃣ 每 格 写 1 个 词 (1 周 学 过 的)",
    "3️⃣ 在 词 旁边 画 一 幅 小 画",
    "4️⃣ 给 每 个 词 贴 一 张 贴纸",
    "5️⃣ 给 你 的 小 书 起 一 个 名字!",
]
for i,line in enumerate(steps):
    tb(s,0.55,2.10+i*0.50,4.30,0.40,line,sz=11,b=True,c=DARK)

# RIGHT — example word list
panel(s,5.05,1.50,4.55,3.50,GOLD)
panel_head(s,5.05,1.50,4.55,GOLD,"📚 词 库 · 选 8 个 词  Word Bank",sz=12)
words=[
    "🌞 太阳   🌙 月亮   🌍 地球",
    "⭐ 星星   🌌 银河   🪐 星球",
    "🚀 火箭   🌕 月球   🔴 火星",
    "🛰️ 太空 站   👽 外星人",
    "🌱 生命   📡 信号   💭 猜想",
]
for i,line in enumerate(words):
    tb(s,5.20,2.10+i*0.45,4.30,0.40,line,sz=11,b=True,c=DARK)

# Bottom tip
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(5.10),Inches(9.20),Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb=ALIEN; tip.line.fill.background()
tb(s,0.55,5.15,9.00,0.30,"💡 完成 后 带 回家 · 给 爸爸 妈妈 读 一 遍!  Take home & read to family!",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"20-30 分钟:\n• 材料: A4 纸 (老师 预 折 更 好) + 贴 纸 + 彩 笔\n• K-2 学生 可以 只 画 + 1 个 字\n• 高 年级 也 可以 选 — 把 8 个 词 都 写 拼音 + 英文")

# ============================================================
# Slide 12 · Project 2 · 知识 竞赛 卡
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🃏 项目 2 · 知识 竞赛 卡  Trivia Card Game",ORANGE)

lvl=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.30),Inches(0.95),Inches(2.30),Inches(0.42))
lvl.fill.solid(); lvl.fill.fore_color.rgb=ORANGE; lvl.line.fill.background()
tb(s,7.30,0.99,2.30,0.34,"🌟 All · 小组 Group",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)

# LEFT — steps
panel(s,0.40,1.50,4.55,3.50,ORANGE)
panel_head(s,0.40,1.50,4.55,ORANGE,"📝 怎么 做  Steps",sz=12)
steps2=[
    "1️⃣ 4-5 人 一 组, 每 组 拿 10 张 卡",
    "2️⃣ 一 面 写 问题, 另 一 面 写 答案",
    "3️⃣ 题目 来 自 这 周 的 学习 内容",
    "4️⃣ 写 完 后 — 和 别 的 组 交换 玩!",
    "5️⃣ 答 对 1 题 = 1 ⭐  最 多 ⭐ 赢!",
]
for i,line in enumerate(steps2):
    tb(s,0.55,2.10+i*0.50,4.30,0.40,line,sz=11,b=True,c=DARK)

# RIGHT — example questions
panel(s,5.05,1.50,4.55,3.50,COSMIC)
panel_head(s,5.05,1.50,4.55,COSMIC,"💡 题目 例子  Example Questions",sz=12)
examples=[
    "❓ 太阳系 里 最 大 的 行星?",
    "❓ 火星 是 什么 颜色?",
    "❓ 北极星 在 哪 个 方向?",
    "❓ Zorp 来 自 哪 里?",
    "❓ 「生命」 的 拼音?",
]
for i,line in enumerate(examples):
    tb(s,5.20,2.10+i*0.45,4.30,0.40,line,sz=11,b=True,c=DARK)
tb(s,5.20,4.50,4.30,0.30,"💬 也 可以 让 队员 选 难度: 简单 / 中等 / 难",sz=9,b=True,c=COSMIC,a=PP_ALIGN.CENTER)

tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(5.10),Inches(9.20),Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb=ORANGE; tip.line.fill.background()
tb(s,0.55,5.15,9.00,0.30,"🎯 每 组 出 题 + 答 题 — 全员 参与!  Everyone contributes & plays!",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30-40 分钟:\n• 材料: 索引 卡 / 厚纸 (每 组 10-15 张) + 笔\n• 出题 时 老师 巡视, 帮 学生 改 错别字\n• 分层: 低 年级 写 简单 配 对 题; 高 年级 写 开放 问 题\n• 玩 法: 每 组 抽 一 张, 队友 一起 答, 答 对 留 着, 答 错 还 回 去")

# ============================================================
# Slide 13 · Project 3 · 我 的 太空 日记 (Grade 3-5)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📔 项目 3 · 我 的 太空 日记  My Space Journal",EARTH)

lvl=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.30),Inches(0.95),Inches(2.30),Inches(0.42))
lvl.fill.solid(); lvl.fill.fore_color.rgb=EARTH; lvl.line.fill.background()
tb(s,7.30,0.99,2.30,0.34,"📖 3-5 · 个人 Individual",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)

# LEFT — page layout
panel(s,0.40,1.50,4.55,3.50,EARTH)
panel_head(s,0.40,1.50,4.55,EARTH,"📝 4 页 · 一 天 一 页  4 Pages · One per Day",sz=12)
pages=[
    ("Day 1","🌞 我 学了 太阳系", GOLD),
    ("Day 2","⭐ 我 认识 了 星座", NEBULA),
    ("Day 3","🚀 我 想象 自己 是 宇航员", MARS),
    ("Day 4","👽 我 设计 了 一个 外星朋友", ALIEN),
]
for i,(d,line,cl) in enumerate(pages):
    y=2.10+i*0.55
    badge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(y),Inches(0.85),Inches(0.40))
    badge.fill.solid(); badge.fill.fore_color.rgb=cl; badge.line.fill.background()
    tb(s,0.55,y+0.06,0.85,0.30,d,sz=10,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.50,y+0.05,3.35,0.35,line,sz=11,b=True,c=DARK)

# RIGHT — sentence frames
panel(s,5.05,1.50,4.55,3.50,COSMIC)
panel_head(s,5.05,1.50,4.55,COSMIC,"✍️ 每 页 写  Each Page Writes",sz=12)
frames=[
    "「我 学 了 ___」",
    "「我 觉得 ___ 最 有趣」",
    "「因为 ___」",
    "「我 还 想 知道 ___」",
    "+ 画 一 幅 小 画!",
]
for i,line in enumerate(frames):
    tb(s,5.20,2.10+i*0.50,4.30,0.40,line,sz=12,b=True,c=COSMIC if i<4 else DARK)
tb(s,5.20,4.65,4.30,0.30,"📷 也 可以 加 照片 / 贴 纸 / 标 题 设计",sz=9,c=GRAY,a=PP_ALIGN.CENTER)

tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(5.10),Inches(9.20),Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb=EARTH; tip.line.fill.background()
tb(s,0.55,5.15,9.00,0.30,"💎 写 完 — 给 朋友 读 一 段, 听 他 给 你 「奖 励 ⭐」!",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30-40 分钟:\n• 材料: 4 页 装订 小 本 (老师 提前 准备) 或 让 学生 自 折\n• 适合 写字 能力 强 的 学生\n• 鼓励 用 句型 + 自由 发挥\n• 写 完 后 可以 自愿 上 来 朗读")

# ============================================================
# Slide 14 · Project 4 · 太空 小 剧场
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎭 项目 4 · 太空 小 剧场  Space Mini Theater",PINK)

lvl=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.30),Inches(0.95),Inches(2.30),Inches(0.42))
lvl.fill.solid(); lvl.fill.fore_color.rgb=PINK; lvl.line.fill.background()
tb(s,7.30,0.99,2.30,0.34,"🎬 All · 小组 Group",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)

# LEFT — scenarios (pick one)
panel(s,0.40,1.50,4.55,3.50,PINK)
panel_head(s,0.40,1.50,4.55,PINK,"🎬 选 一 个 剧本  Pick a Scenario",sz=12)
scenarios=[
    ("A","🚀 宇航员 第 一 次 到 火星, 见 到 Zorp"),
    ("B","✉️ 地球 小朋友 收 到 Zorp 的 信"),
    ("C","⭐ 古人 看 星 空, 发现 新 星座"),
    ("D","📡 科学家 收 到 太空 信号 — 是 谁?"),
]
for i,(letter,line) in enumerate(scenarios):
    y=2.10+i*0.60
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y),Inches(0.50),Inches(0.50))
    badge.fill.solid(); badge.fill.fore_color.rgb=PINK; badge.line.fill.background()
    tb(s,0.55,y+0.10,0.50,0.30,letter,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.15,y+0.10,3.65,0.45,line,sz=11,b=True,c=DARK)

# RIGHT — steps
panel(s,5.05,1.50,4.55,3.50,COSMIC)
panel_head(s,5.05,1.50,4.55,COSMIC,"📝 怎么 排  Steps",sz=12)
steps4=[
    "1️⃣ 3-4 人 一 组, 选 一 个 剧本",
    "2️⃣ 每 人 选 一 个 角色",
    "3️⃣ 想 3-5 句 中文 台词",
    "4️⃣ 用 纸 + 彩 笔 做 简单 道具",
    "5️⃣ 排 练 — 然后 上 台 表演!",
    "6️⃣ 长度: 1-2 分钟",
]
for i,line in enumerate(steps4):
    tb(s,5.20,2.10+i*0.45,4.30,0.40,line,sz=11,b=True,c=DARK)

tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(5.10),Inches(9.20),Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb=PINK; tip.line.fill.background()
tb(s,0.55,5.15,9.00,0.30,"🎤 表演 后 — 全 班 鼓 掌 + 拍 照 留 念!  Bow & take a class photo!",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30-45 分钟:\n• 材料: 纸 + 彩 笔 (做 头 饰) + 简单 道具\n• 老师 不 要 写 台词 — 让 学生 自由 发挥\n• 鼓励 多 用 中文 句子 (但 也 可以 加 英文 单词)\n• 拍 视频 发给 家长 (征 得 同 意 后)")

# ============================================================
# Slide 15 · 分享 + 毕业 + 结业 证书
# ============================================================
s=ns(); bg(s,NIGHT)

# Stars
random.seed(99)
for _ in range(40):
    x = random.uniform(0.2, 9.7); y = random.uniform(0.2, 5.4); sz = random.choice([0.06,0.10])
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()

tb(s,0.3,0.25,9.4,0.55,"🎓 毕业 啦!  Congrats, Space Explorers!",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)

# Certificate card
cert=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.00),Inches(1.10),Inches(8.00),Inches(2.65))
cert.fill.solid(); cert.fill.fore_color.rgb=CREAM
cert.line.color.rgb=STAR; cert.line.width=Pt(4)
tb(s,1.00,1.25,8.00,0.40,"🏅  太空 探险家 结业 证书",sz=18,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,1.00,1.70,8.00,0.32,"Certificate of Space Exploration",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1.30,2.15,7.40,0.36,"恭喜  ____________________________",sz=14,b=True,c=DARK)
tb(s,1.30,2.55,7.40,0.32,"完成 5 天 太空 之 旅 — 学习 了 太阳系 · 星座 · 航天 · 外星人!",sz=11,c=DARK)
tb(s,1.30,2.92,7.40,0.30,"You completed a 5-day space journey — Solar System · Stars · Spaceflight · Aliens!",sz=9,c=GRAY)
tb(s,1.30,3.35,7.40,0.32,"日期 Date: __________   老师 Teacher: __________",sz=11,c=DARK)

# Share box
share=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.50),Inches(4.00),Inches(9.00),Inches(1.20))
share.fill.solid(); share.fill.fore_color.rgb=COSMIC
share.line.color.rgb=STAR; share.line.width=Pt(2.5)
tb(s,0.65,4.08,8.70,0.32,"🎤 分享 时间  Share Time",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.65,4.45,8.70,0.32,"「我 最 喜欢 的 是 ___ 因为 ___」",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.65,4.82,8.70,0.30,"My favorite was ___ because ___",sz=10,c=LGRAY,a=PP_ALIGN.CENTER)

tb(s,0.5,5.25,9.0,0.30,"👋 再 见 — 继续 仰望 星 空!  See you · keep looking up!",sz=11,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"最 后 10-15 分钟:\n• 每 个 学生 / 小组 分享 1 分钟\n• 颁 「太空 探险家」证书 (老师 提前 打印)\n• 全 班 合 影\n• 鼓励 学生 把 项目 带 回 家")

out = os.path.join(os.path.dirname(__file__), "day5_review_projects.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
