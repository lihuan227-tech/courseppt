#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
仰望星空 Looking Up at the Stars · Day 1: 太阳系和宇宙奥秘  Solar System & Cosmic Wonders
Picture book anchor: 神奇校车 · 迷失太阳系  The Magic School Bus — Lost in the Solar System
3 sessions × ~50 min · K-5
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# === Palette: Cosmic Night ===
NIGHT  = RGBColor(0x0D,0x1B,0x3E)   # deep night sky
COSMIC = RGBColor(0x6A,0x1B,0x9A)   # nebula purple
STAR   = RGBColor(0xF5,0xC2,0x42)   # golden star
GOLD   = RGBColor(0xFF,0xB7,0x00)   # bright gold
EARTH  = RGBColor(0x1E,0x88,0xE5)   # earth blue
GREEN  = RGBColor(0x2E,0x7D,0x32)   # earth green
MOON   = RGBColor(0xB0,0xBE,0xC5)   # silver moon
MARS   = RGBColor(0xD8,0x43,0x15)   # mars red
SUN    = RGBColor(0xFF,0x8F,0x00)   # bright sun amber
NEBULA = RGBColor(0x7B,0x1F,0xA2)
SKY    = RGBColor(0x42,0xA5,0xF5)
PINK   = RGBColor(0xEC,0x40,0x7A)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)

# === Helpers ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=IMGBG; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=STAR,a=PP_ALIGN.CENTER)
    # Star sprinkles
    for x,y in [(0.8,4.7),(1.8,4.5),(7.8,4.5),(8.6,4.7),(2.0,1.0),(8.0,1.0)]:
        d=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(x),Inches(y),Inches(0.35),Inches(0.35))
        d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
    return s
def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mid_x=x+size/2; mid_y=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mid_x-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(mid_y-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

n=0

# ============================================================
# 1. COVER
# ============================================================
s=ns(); bg(s,NIGHT)
# Star sprinkles in night sky
import random
random.seed(7)
for _ in range(45):
    x=random.uniform(0.3,9.7); y=random.uniform(0.3,5.3); sz=random.uniform(0.07,0.18)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
tb(s,0.5,0.4,9,0.6,"🌌 仰望星空  Looking Up at the Stars",sz=28,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.0,9,0.45,"Day 1 · 太阳系和宇宙奥秘  Solar System & Cosmic Wonders",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# IMAGE PLACEHOLDER — replace with real cover image (e.g., 太阳系 / 绘本封面)
ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.5),Inches(1.75),Inches(5.0),Inches(2.85))
ph.fill.solid(); ph.fill.fore_color.rgb=IMGBG; ph.line.color.rgb=STAR; ph.line.width=Pt(3)
tb(s,2.5,2.55,5.0,0.7,"🖼️",sz=60,a=PP_ALIGN.CENTER)
tb(s,2.5,3.30,5.0,0.36,"图片 位置  Image Placeholder",sz=14,b=True,c=NIGHT,a=PP_ALIGN.CENTER)
tb(s,2.5,3.70,5.0,0.28,"老师 后期 放 真实 图片",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,2.5,4.00,5.0,0.28,"Teacher: insert real image later",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.5,4.75,9,0.40,"📖 绘本: 神奇校车 · 迷失太阳系",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,5.15,9,0.28,"Picture book: The Magic School Bus — Lost in the Solar System",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 1 开场 (1 分钟):\n• 「今天我们坐 神奇校车 — 去 太阳系 旅行!」\n• 介绍绘本: 神奇校车 · 迷失太阳系\n• 全班 4 队 (红/蓝/绿/黄), 每队 拿 一颗 行星 名牌")

# ============================================================
# 2. SESSION 1 DIVIDER
# ============================================================
s=div("Session 1  上午 11:00–11:45","☀️ 故事课 · 太阳系大冒险  ·  45 min",COSMIC,"🚀"); n+=1; pn(s,n)

# ============================================================
# 3. LEARNING GOALS
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎯 今天的学习目标  Today's Learning Goals",NIGHT)
tb(s,0.4,0.85,9.2,0.30,"上完这节课, 你会……  By the end, you'll be able to…",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
goals=[
    ("1","☀️","认识 太阳系 — 太阳、月亮、地球、其他 行星",COSMIC),
    ("2","🌍","知道 地球 只是 宇宙 中 的 一 小 部分",EARTH),
    ("3","🔭","会 比较 行星 — 大小、冷热、远近",MARS),
    ("4","💡","对 宇宙 探索 充满 好奇 — 大胆 提问!",STAR),
]
for i,(num,em,text,cl) in enumerate(goals):
    y=1.30+i*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.18),Inches(0.50),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.22,0.50,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.18,0.60,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.90,y+0.20,7.6,0.55,text,sz=14,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"1-2 分钟 — 学习目标:\n• 老师 念 4 个 目标\n• 提示: 「今天 我们 一起 去 太空 看看!」")

# ============================================================
# 4. HOOK — 你 抬头 看过 天空 吗?
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔭 你 抬头 看过 天空 吗?  Have You Looked Up at the Sky?",STAR)
tb(s,0.4,0.85,9.2,0.36,"白天、晚上 — 天上 有 什么? 想 一 想!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Daytime, nighttime — what's up in the sky? Think about it!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
items=[
    ("☀️","太阳","Sun","白天 看到 — 又亮 又暖",SUN),
    ("🌙","月亮","Moon","晚上 看到 — 有 时 圆 有 时 弯",MOON),
    ("⭐","星星","Stars","晚上 闪闪 亮 — 一闪一闪",STAR),
    ("☁️","云","Clouds","飘 来 飘 去",SKY),
]
for i,(em,cn,en,d,cl) in enumerate(items):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.22),Inches(2.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.12,1.0,em,sz=66,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,2.12,0.36,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.12,2.12,0.22,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.40,2.02,0.50,d,sz=10,c=DARK,a=PP_ALIGN.CENTER)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.25),Inches(9.2),Inches(1.15))
prompt.fill.solid(); prompt.fill.fore_color.rgb=NIGHT; prompt.line.color.rgb=STAR; prompt.line.width=Pt(2.5)
tb(s,0.55,4.35,9.0,0.30,"🙋 你 还 看到 过 什么? 举手 说说看!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.68,9.0,0.30,"What else have you seen up there? Share!",sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,4.98,9.0,0.30,"💬 「我 看到 过 ___ 在 天上。」",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🔥 HOOK · 3-4 分钟:\n• 老师 问: 「你 最近 抬头 看过 天空 吗? 看到 什么?」\n• 收集 3-5 个 答案 — 写 在 黑板\n• 引导: 「天上 这么多 东西 — 它们 都 在 一个 大 地方 叫 太阳系!」")

# ============================================================
# 5. PICTURE BOOK INTRO — 神奇校车
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📖 绘本: 神奇校车 · 迷失太阳系",COSMIC)
tb(s,0.4,0.85,9.2,0.40,"Ms. Frizzle 和 全班 — 坐 校车 飞 上 太空!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"Ms. Frizzle and her class — riding the bus into space!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.5,1.65,4.5,2.95,"📚 绘本 封面 / Book cover")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(2.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=COSMIC; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=COSMIC; head.line.fill.background()
tb(s,5.40,1.72,4.10,0.40,"📚 故事 简介  About the Story",sz=13,b=True,c=WHITE)
tb(s,5.40,2.28,4.10,0.34,"作者: Joanna Cole",sz=12,b=True,c=DARK)
tb(s,5.40,2.62,4.10,0.30,"Author · Magic School Bus 系列",sz=9,c=GRAY)
tb(s,5.40,2.98,4.10,0.34,"📍 一群 小朋友 跟 老师 Ms. Frizzle 一起 坐 校车",sz=11,b=True,c=DARK)
tb(s,5.40,3.32,4.10,0.34,"   飞 上 太空, 一站 一站 看 行星!",sz=11,b=True,c=DARK)
tb(s,5.40,3.70,4.10,0.34,"⭐ 他们 还 不小心 「迷路」 了!",sz=11,b=True,c=COSMIC)
tb(s,5.40,4.10,4.10,0.34,"They got lost in space — what now?",sz=9,c=GRAY)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — 介绍绘本:\n• 展示 绘本 封面 — 让 学生 猜 这 是 一辆 什么样 的 校车?\n• 介绍 Ms. Frizzle (老师) 和 学生\n• 引导: 「他们 要 去 太阳系 旅行 — 我们 一起 跟着 看!」")

# ============================================================
# 5b. 绘本前 · Pre-reading predictions  (Think-Pair-Share)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔮 翻 开 之前 — 想 一 想!  Before We Read",COSMIC)
tb(s,0.4,0.85,9.2,0.34,"光 看 封面 — 你 能 猜 到 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Just from the cover — what can you guess?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
preds=[
    ("🚌","校车 真的 会 飞 吗?","Can a bus really fly?",COSMIC),
    ("🪐","第 一站 去 哪个 行星?","Which planet first?",STAR),
    ("😱","如果 「迷路」 — 怎么办?","What if they get LOST?",MARS),
]
for i,(em,q,en,cl) in enumerate(preds):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.0,em,sz=66,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,2.85,0.50,q,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.30,2.85,0.40,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.80,2.75,0.55,"💡 没有 标准 答案 — 大胆 猜!",sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
tps=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.55),Inches(9.2),Inches(0.95))
tps.fill.solid(); tps.fill.fore_color.rgb=NIGHT; tps.line.color.rgb=STAR; tps.line.width=Pt(2.5)
tb(s,0.55,4.62,9.0,0.30,"👥 Think-Pair-Share: 跟 同桌 说 你 的 猜 想 (1 分钟)",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.26,"Turn to a partner — share your guess!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.24,"💬 「我 猜 ___ 因为 ___」",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🔮 PRE-READING · 3 分钟:\n• 让 学生 先 猜 — 激发 好奇 心\n• 没有 对错 — 收集 想法\n• 写 1-2 个 在 黑板 上 — 读完 验证\n• K-1 也能 说 — 「我 猜 太阳!」 就够")

# ============================================================
# 5c. 听 绘本 时 · During-reading observation checklist
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"👀 听 故事 时 — 你 的 任务!  While You Listen",STAR)
tb(s,0.4,0.85,9.2,0.34,"3 个 任务 — 一边 听 一边 留心!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"3 missions — listen + observe at the same time!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
obs=[
    ("1","🪐","他们 去 了 哪些 行星?","Which planets did they visit?",NEBULA),
    ("2","😱","哪个 行星 最 让 他们 害怕?","Which planet scared them most?",MARS),
    ("3","💡","他们 怎么 回到 地球?","How did they get back to Earth?",EARTH),
]
for i,(num,em,cn,en,cl) in enumerate(obs):
    y=1.55+i*1.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.95))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.20),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.27,0.55,0.40,num,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.18,0.80,0.55,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,2.30,y+0.14,7.0,0.40,cn,sz=15,b=True,c=cl)
    tb(s,2.30,y+0.54,7.0,0.30,en,sz=10,c=GRAY)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.60))
tip.fill.solid(); tip.fill.fore_color.rgb=STAR; tip.line.fill.background()
tb(s,0.55,4.93,9.0,0.30,"👂 用 心 听! 念 完 之后 — 我们 一起 讨论",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,5.23,9.0,0.22,"Listen well — we'll discuss after the story.",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"👀 DURING-READING · 1 分钟 setup + 8-10 分钟 念书:\n• 念 3 个 任务 — 学生 知道 听 什么\n• 老师 加 手势 (耳朵=听, 眼睛=看, 脑子=想)\n• 然后 念 绘本 (选 重点 段落)\n• 念 完 → 下一页 讨论")

# ============================================================
# 5d. 绘本 后 · Post-reading discussion (6 Qs)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"💭 听 完 故事 — 一起 讨论!  After We Read",COSMIC)
tb(s,0.4,0.85,9.2,0.30,"选 1-2 个 问题 — 全班 / 小组 / Think-Pair-Share",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.15,9.2,0.22,"Pick 1-2 questions — class / group / Think-Pair-Share",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("🚌","校车 飞 上 太空 — 神奇 吗?","Was the flying bus magical?"),
    ("🪐","哪个 行星 最 危险? 为什么?","Most dangerous planet? Why?"),
    ("🦸","如果 你 在 校车 上 — 你 会 怎么 帮?","If YOU were there — how would you help?"),
    ("😨","学生 们 害怕 吗? 你 害怕 吗?","Were they scared? Would you be?"),
    ("👩‍🏫","Ms. Frizzle 厉害 在 哪里?","What makes Ms. Frizzle awesome?"),
    ("❓","看完 后 — 你 还 想 知道 什么?","After reading — what do YOU wonder?"),
]
for i,(em,cn,en) in enumerate(qs):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.45+row*1.80
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.65))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=COSMIC; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.42),Inches(0.42))
    nb.fill.solid(); nb.fill.fore_color.rgb=COSMIC; nb.line.fill.background()
    tb(s,x+0.10,y+0.14,0.42,0.34,str(i+1),sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.60,y+0.05,0.55,0.50,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,y+0.60,2.70,0.65,cn,sz=11,b=True,c=DARK)
    tb(s,x+0.15,y+1.25,2.70,0.32,en,sz=8,c=GRAY)
n+=1; pn(s,n)
notes(s,"💭 POST-READING · 5-7 分钟:\n• 选 2-3 题 — 全班 / Think-Pair-Share\n• Q3 + Q6 最 启发 想象力\n• 答错 没关系 — 收集 想法, 不给 标准答案\n• 高年级 完整 句子 / 低年级 几个 字 也行")

# ============================================================
# 6. SOLAR SYSTEM OVERVIEW — 8 planets
# ============================================================
s=ns(); bg(s,NIGHT); hb(s,"🌌 太阳系 · 我们 的 家  Our Solar System",STAR,t=0.15)
tb(s,0.4,0.85,9.2,0.30,"☀️ 1 个 太阳 + 8 个 行星 — 一个 大家庭!",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"1 sun + 8 planets — one big family!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
planets=[
    ("☀️","太阳","Sun",SUN),
    ("☿","水星","Mercury",LGRAY),
    ("♀","金星",  "Venus",GOLD),
    ("🌍","地球","Earth",EARTH),
    ("🔴","火星","Mars",MARS),
    ("🪐","木星","Jupiter",MARS),
    ("💍","土星","Saturn",STAR),
    ("🔵","天王星","Uranus",SKY),
    ("🔷","海王星","Neptune",EARTH),
]
for i,(em,cn,en,cl) in enumerate(planets):
    col=i%5; row=i//5
    x=0.40+col*1.88; y=1.55+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.20),Inches(y+0.10),Inches(1.40),Inches(1.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=cl; sh.line.color.rgb=STAR; sh.line.width=Pt(2)
    tb(s,x+0.20,y+0.35,1.40,0.80,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+1.55,1.60,0.28,cn,sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+1.80,1.60,0.20,en,sz=8,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,5.25,9.2,0.30,"📏 离 太阳: 水星 最近 · 海王星 最远",sz=11,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5-6 分钟 — 太阳系大家庭:\n• 念 9 个 名字 (太阳 + 8 行星) — 全班 跟读\n• 记忆 口诀: 水金地火, 木土天海 (顺序!)\n• 互动: 让 4 队 各 演 一个 行星 — 围着 老师(太阳) 转一圈")

# ============================================================
# 7. PLANET COMPARISONS — size, temperature, distance
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📊 行星 比一比!  Let's Compare Planets!",MARS)
tb(s,0.4,0.85,9.2,0.30,"大小、冷热、远近 — 看看 谁 最 ___?",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Size · temperature · distance — who's most ___?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
cmps=[
    ("🏆","最大 的 行星","Biggest","🪐 木星 Jupiter","可以 装 下 1300 个 地球!",NEBULA),
    ("🥶","最冷 的 行星","Coldest","🔷 海王星 Neptune","-200°C — 比 冰箱 还 冷!",SKY),
    ("🥵","最热 的 行星","Hottest","♀ 金星 Venus","460°C — 比 烤箱 还 热!",MARS),
    ("🏠","我们 的 家","Our home","🌍 地球 Earth","唯一 有 生命 — 有 水, 有 空气!",EARTH),
]
for i,(em,cn,en,winner,fact,cl) in enumerate(cmps):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.15,0.95,0.85,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,x+1.10,y+0.12,3.35,0.36,cn,sz=14,b=True,c=cl)
    tb(s,x+1.10,y+0.48,3.35,0.26,en,sz=9,c=GRAY)
    tb(s,x+1.10,y+0.78,3.35,0.36,winner,sz=14,b=True,c=DARK)
    tb(s,x+1.10,y+1.15,3.35,0.55,fact,sz=10,c=DARK)
n+=1; pn(s,n)
notes(s,"4 分钟 — 行星 比较:\n• 一对 一对 念 — 让 学生 张大 眼睛 「哇!」\n• 重点: 地球 是 唯一 有 生命 的 — 因为 有 水 + 空气 + 合适 的 温度\n• 引导: 「所以 我们 要 保护 地球!」")

# ============================================================
# 7b. TPR 行星 表演 · Be a Planet (whole-class movement)
# ============================================================
s=ns(); bg(s,NIGHT); hb(s,"🎭 我 是 行星! 来 演 一下!  Be a Planet!",STAR)
tb(s,0.4,0.85,9.2,0.34,"全班 起 立! 老师 是 太阳, 你们 是 行星!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Everyone stand up! Teacher = Sun, you = Planets!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
actions=[
    ("☀️","太阳","Sun","站 中间 不动 — 双手 「发光」",SUN),
    ("☿","水星","Mercury","跑 最 快 — 围 老师 转 小圈",LGRAY),
    ("🌍","地球","Earth","稳稳 走 + 一边 自转",EARTH),
    ("🪐","木星","Jupiter","巨大 步伐 — 大 大 + 慢",NEBULA),
    ("💍","土星","Saturn","张 开 双手 (大 环!) 转",STAR),
    ("🔷","海王星","Neptune","慢 慢 转 大 圈 — 远 远",SKY),
]
for i,(em,cn,en,action,cl) in enumerate(actions):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.55+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.10,0.80,0.85,em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x+1.00,y+0.10,1.85,0.36,cn,sz=14,b=True,c=cl)
    tb(s,x+1.00,y+0.46,1.85,0.26,en,sz=9,c=GRAY)
    tb(s,x+0.12,y+0.85,2.75,0.50,action,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.78,9.2,0.30,"🎵 升级 玩法: 老师 喊 「水金地火 木土天海」 — 学生 顺序 转!",sz=11,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,5.10,9.2,0.24,"Bonus: teacher chants the order — students rotate in sequence!",sz=8,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎭 TPR · 5-6 分钟 — 行星 表演:\n• 全班 起立 — 站 大 圆 圈\n• 老师 站 中间 当 太阳 — 双手 张开 「发光」\n• 老师 喊 一个 行星 — 学生 全部 演 那个 动作 (10 秒)\n• 玩 2 轮 — 第二 轮 加 顺序 挑战\n• 出汗 + 笑 = 学进 大脑! K-5 都 喜欢\n• 课堂管理: 提前 划 圈, 不撞人")

# ============================================================
# 7c. 快问快答 · Rapid Fire Quiz (team buzz)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🚀 快 问 快 答!  Rapid Fire Quiz!",MARS)
tb(s,0.4,0.85,9.2,0.34,"4 队 抢答 — 答 对 加 1 分! 答 错 不 扣 分!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"4 teams buzz in — correct = +1 · wrong = no penalty",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
quiz=[
    ("1","🌍","哪 个 是 我们 的 家?","→ 地球",EARTH),
    ("2","🔷","最 远 的 行星?","→ 海王星",SKY),
    ("3","♀","最 热 是 哪 个?","→ 金星 (460°C)",MARS),
    ("4","🪐","最 大 的 行星?","→ 木星",NEBULA),
    ("5","💍","哪 颗 有 大 环?","→ 土星",STAR),
    ("6","☀️","太阳 是 行星 吗?","→ 不是! 是 恒星",SUN),
]
for i,(num,em,q,ans,cl) in enumerate(quiz):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.55+row*1.70
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.42),Inches(0.42))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+0.10,y+0.14,0.42,0.34,num,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.60,y+0.08,0.65,0.45,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,y+0.58,2.70,0.40,q,sz=13,b=True,c=DARK)
    tb(s,x+0.15,y+1.05,2.70,0.40,ans,sz=12,b=True,c=cl)
n+=1; pn(s,n)
notes(s,"🚀 QUIZ · 4-5 分钟:\n• 4 队 — 举手 / 拍 桌子 抢答\n• 老师 念 题 — 第一个 抢到 的 答\n• 答 错 — 下一队 接\n• 留 时间 强调 Q6: 「太阳 不是 行星 — 是 恒星! 自己 会 发光」")

# ============================================================
# 8. SESSION 1 SHARE / WRAP
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎤 一起 想 一 想  Discuss Together",COSMIC)
tb(s,0.4,0.85,9.2,0.32,"听完 故事 + 看完 行星 — 你 觉得……",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"After the story + the planets — what do YOU think?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("🤔","你 最 喜欢 哪个 行星? 为什么?","Which planet do you like best? Why?",STAR),
    ("🌍","为什么 地球 是 我们 的 家?","Why is Earth our home?",EARTH),
    ("🚀","如果 你 也 能 坐 神奇校车, 你 想 去 哪里?","If YOU could ride the bus — where to?",COSMIC),
]
for i,(em,q,en,cl) in enumerate(qs):
    y=1.55+i*1.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.00))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,0.55,y+0.22,0.80,0.60,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.18,8.0,0.40,q,sz=15,b=True,c=cl)
    tb(s,1.45,y+0.58,8.0,0.30,en,sz=10,c=GRAY)
n+=1; pn(s,n)
notes(s,"5 分钟 — 讨论 + 收 尾:\n• 选 1-2 个 问题 全班 讨论\n• 让 3-4 个 学生 分享 (轮流 — 高低 年级 都 试)\n• 引出 Session 2: 「下午 — 我们 来 学 太空 词语!」")

# ============================================================
# 9. SESSION 2 DIVIDER
# ============================================================
s=div("Session 2  下午 2:00–2:45","📚 词汇课 · 我会认 + 我会写  ·  45 min",EARTH,"📖"); n+=1; pn(s,n)

# ============================================================
# 10. REVIEW — 早上学了什么?
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔁 早上 学了 什么?  Morning Review",EARTH)
tb(s,0.4,0.85,9.2,0.32,"想 一 想 — 早上 我们 一起 看了 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"What did we explore this morning?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
items=[
    ("🚌","神奇 校车","Magic School Bus","坐 校车 去 太空!",COSMIC),
    ("☀️","太阳 + 8 行星","Sun + 8 planets","水金地火 木土天海",STAR),
    ("🌍","地球","Earth","唯一 有 生命 的 家",EARTH),
    ("📊","行星 比较","Compare planets","最大 最冷 最热!",MARS),
]
for i,(em,cn,en,d,cl) in enumerate(items):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.20,1.00,1.20,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.18,3.25,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.20,y+0.58,3.25,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.92,3.25,0.70,d,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"2-3 分钟 — 复习:\n• 让 学生 轮流 说 一句 — 「我 记得 ___」\n• 4 个 主要 点 全 cover 后 — 进入 词汇")

# ============================================================
# 11-15. 我会认 — 5 vocabulary words (one per slide, matching other-unit format)
# ============================================================
read_words=[
    ("☀️","太阳","tài yáng","Sun",SUN,
        "太阳 让 大地 变 暖。",
        "📷 太阳 / 蓝天 / 暖光"),
    ("🌙","月亮","yuè liang","Moon",MOON_C := RGBColor(0xB0,0xBE,0xC5),
        "晚上 月亮 在 天上 发光。",
        "📷 月亮 / 夜空 / 圆月"),
    ("🌍","地球","dì qiú","Earth",EARTH,
        "地球 是 我们 的 家。",
        "📷 地球 / 蓝色 星球 / 大陆"),
    ("🪐","星球","xīng qiú","Planet",NEBULA,
        "太阳系 有 8 个 星球。",
        "📷 行星 / 土星 / 木星"),
    ("🌌","宇宙","yǔ zhòu","Universe",COSMIC,
        "宇宙 里 有 好多 星星。",
        "📷 银河 / 星云 / 深空"),
]
for em,cn,py,en,c,sent,img_label in read_words:
    s=ns(); bg(s,CREAM); hb(s,f"👀 我会认 · {cn}  I Can Read",c)
    # Left: big character card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.6))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.fill.background()
    tb(s,0.5,1.05,4.3,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,0.5,1.75,4.3,1.0,cn,sz=66 if len(cn)==2 else 72,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,f"{py}  ·  {en}",sz=18,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.25,4.3,0.34,"👉 跟我读!  Read after me!",sz=13,b=True,c=c,a=PP_ALIGN.CENTER)
    # Right: image placeholder
    ib_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.3),Inches(1.0),Inches(4.4),Inches(2.6))
    ib_box.fill.solid(); ib_box.fill.fore_color.rgb=IMGBG; ib_box.line.color.rgb=c; ib_box.line.width=Pt(2)
    tb(s,5.3,1.80,4.4,0.6,"🖼️",sz=44,a=PP_ALIGN.CENTER)
    tb(s,5.3,2.50,4.4,0.4,img_label,sz=12,c=LGRAY,a=PP_ALIGN.CENTER)
    tb(s,5.3,2.95,4.4,0.30,"图片 位置 · Image placeholder",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
    # Bottom: example sentence
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.30))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=c; sh2.line.width=Pt(2.5)
    tb(s,0.6,3.95,2.0,0.40,"📌 例句  Example",sz=14,b=True,c=c)
    tb(s,0.6,4.40,8.8,0.55,sent,sz=22,b=True,c=DARK)
    tb(s,0.4,5.25,9.2,0.28,"💬 「我 认识 ___」  · I know the word ___",sz=11,b=True,c=c,a=PP_ALIGN.CENTER)
    n+=1; pn(s,n)
    notes(s,f"3 分钟 — {cn}:\n• 老师 指 字, 全班 齐读 3 遍 (慢 → 快 → 大声)\n• 看 图: 「这 是 ___, 你 见过 吗?」\n• 读 例句, 学生 跟读\n• 抽 1-2 个 学生 用「{cn}」造 一 个 新 句子\n• 写 到 黑板 上 — 让 学生 在 空中 跟着 写 一遍")

# ============================================================
# 11b. 词汇 闪卡 抢答 · Flash Card Challenge
# ============================================================
s=ns(); bg(s,NIGHT); hb(s,"⚡ 词汇 闪卡 抢答!  Flash Card Challenge!",STAR)
tb(s,0.4,0.85,9.2,0.34,"老师 出 词 — 哪 队 第一 大声 读 + 翻译?",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Teacher flashes a word — first team to read + translate wins!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
rules=[
    ("1️⃣","老师 举 一张 闪卡","Teacher holds up a card"),
    ("2️⃣","第一 队 抢答 — 读 + 翻译","First team — read + translate"),
    ("3️⃣","对 = +1 分, 错 = 下 一队","Right = +1, wrong = next team"),
    ("4️⃣","用 这个 词 造一 句 = +2 分","Make a sentence = +2 bonus"),
]
for i,(num,cn,en) in enumerate(rules):
    y=1.55+i*0.80
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.70))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=STAR; sh.line.width=Pt(2)
    tb(s,0.55,y+0.14,0.6,0.42,num,sz=18,b=True,c=STAR)
    tb(s,1.25,y+0.08,4.7,0.32,cn,sz=14,b=True,c=DARK)
    tb(s,1.25,y+0.42,4.7,0.24,en,sz=9,c=GRAY)
    tb(s,5.95,y+0.20,3.5,0.32,"⚡ 抢 答!",sz=14,b=True,c=MARS,a=PP_ALIGN.CENTER)
tb(s,0.4,4.95,9.2,0.30,"🏆 总冠军 队 — Hi-five + 团队 欢呼!",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"⚡ FLASH CARD · 5-6 分钟:\n• 用 我会认 5 张 卡 (太阳/月亮/地球/星球/宇宙)\n• 玩 2 轮 — 第二 轮 更 快\n• 鼓励 K-1 用 一 个 词, 高年级 完整 翻译\n• 总冠军 — Hi-five + 集体 欢呼")

# ============================================================
# 11c. 造 句 大作战 · Sentence Building (Think-Pair-Share)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"💬 造 句 大作战!  Build a Sentence!",NEBULA)
tb(s,0.4,0.85,9.2,0.34,"用 句型 + 词 — 自己 造 一 句 太空 话!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Use a frame + word — make YOUR space sentence!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
frames=[
    ("🚀","我 想 去 ___ 行星 — 因为 ___","I want to visit ___ planet — because ___"),
    ("📏","___ 离 太阳 最 远","___ is farthest from the sun"),
    ("🏠","___ 是 我们 的 家","___ is our home"),
]
for i,(em,cn,en) in enumerate(frames):
    y=1.55+i*1.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.95))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=NEBULA; sh.line.width=Pt(2.5)
    tb(s,0.55,y+0.18,0.80,0.60,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,1.50,y+0.12,8.0,0.40,cn,sz=17,b=True,c=NEBULA)
    tb(s,1.50,y+0.55,8.0,0.30,en,sz=10,c=GRAY)
bottom=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.65))
bottom.fill.solid(); bottom.fill.fore_color.rgb=NEBULA; bottom.line.fill.background()
tb(s,0.55,4.92,9.0,0.30,"👥 Think-Pair-Share — 同桌 互 说 + 抽 3-4 个 上 台 分享",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.22,"Think — Pair — Share — then a few share with the class!",sz=9,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"💬 SENTENCE BUILDING · 6 分钟:\n• 1 分钟 — 念 3 个 句型\n• 2 分钟 — Turn & Talk (同桌)\n• 3 分钟 — 抽 3-4 个 学生 上 台 — 大声 说\n• 老师 写 在 黑板 上 — 鼓掌\n• 高年级 编 完整 句子, 低年级 填 1-2 个 词 就行")

# ============================================================
# 12-14. 我会写 — 太阳 (2 chars), 地球 (2 chars), 宇宙 (2 chars)
# ============================================================
def write_slide(emoji,word_cn,word_en,chars,color):
    s=ns(); bg(s,CREAM); hb(s,f"✏️ 我会写 · {word_cn}  I Can Write · {word_en}",color)
    tb(s,0.4,0.85,9.2,0.36,f"{emoji} 一起来写「{word_cn}」!",sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.25,9.2,0.26,f"Practice writing {word_cn} ({word_en})",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    if len(chars)==2:
        tianzi(s,0.55,1.65,2.20,chars[0][0],color,pinyin=chars[0][1],char_sz=120)
        tianzi(s,2.95,1.65,2.20,chars[1][0],color,pinyin=chars[1][1],char_sz=120)
    else:
        tianzi(s,1.30,1.65,2.95,chars[0][0],color,pinyin=chars[0][1],char_sz=160)
    # Right panel — how to write
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.45,1.72,4.10,0.40,"✏️ 怎么写  How to Write",sz=13,b=True,c=WHITE)
    for i,(ch,py,hint_cn,hint_en) in enumerate(chars):
        y=2.30+i*0.95
        tb(s,5.45,y,4.10,0.35,f"📐「{ch}」 — {py}",sz=13,b=True,c=DARK)
        tb(s,5.45,y+0.32,4.10,0.30,hint_cn,sz=10,b=True,c=color)
        tb(s,5.45,y+0.58,4.10,0.26,hint_en,sz=9,c=GRAY)
    # Bottom — practice 3 times + sentence frame
    pf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
    pf.fill.solid(); pf.fill.fore_color.rgb=WARM; pf.line.color.rgb=color; pf.line.width=Pt(2)
    tb(s,0.55,4.72,9.0,0.32,f"📝 在 田字格 里 写 3 遍  Practice 3 times in grid paper",sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.55,5.08,9.0,0.32,f"💬 「我 会 写「{word_cn}」!」",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    return s

s=write_slide("☀️","太阳","Sun",[
    ("太","tài","4 笔 — 「大」+ 一点","Like 大 with a dot below"),
    ("阳","yáng","6 笔 — 「阝」+「日」","Left part + sun (日)"),
],SUN); n+=1; pn(s,n)
notes(s,"5-6 分钟 — 写 太阳:\n• 演示 笔顺 — 学生 跟着 空写\n• 记忆: 「太」 = 大 + 一点 (太大!)\n• 「阳」 有 「日」 — 太阳 就 是 日!\n• 田字格 写 3 遍")

s=write_slide("🌍","地球","Earth",[
    ("地","dì","6 笔 — 「土」+「也」","Earth side + 也"),
    ("球","qiú","11 笔 — 「王」+「求」","King radical + 求"),
],EARTH); n+=1; pn(s,n)
notes(s,"5-6 分钟 — 写 地球:\n• 「地」 有 「土」 — 大地 = 泥土\n• 「球」 有 「王」 — 像 一个 大 王 一样 圆\n• 田字格 写 3 遍")

s=write_slide("🌌","宇宙","Universe",[
    ("宇","yǔ","6 笔 — 「宀」+「于」","Roof radical + 于"),
    ("宙","zhòu","8 笔 — 「宀」+「由」","Roof radical + 由"),
],COSMIC); n+=1; pn(s,n)
notes(s,"5-6 分钟 — 写 宇宙:\n• 「宇宙」 都 有 「宀」 (宝盖头) — 像 一个 大 屋顶 罩着 一切\n• 「宇」 + 「宙」 = 所有 空间 + 所有 时间\n• 田字格 写 3 遍")

# ============================================================
# 16. SESSION 3 DIVIDER
# ============================================================
s=div("Session 3  下午 3:00–4:30","🛠️ 项目课 + Day 1 Booklet  ·  90 min",MARS,"🚀"); n+=1; pn(s,n)

# ============================================================
# 17. PROJECTS OVERVIEW — 3 projects
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🛠️ 3 个 项目  3 Projects",MARS)
tb(s,0.4,0.85,9.2,0.32,"选 一个 你 最 喜欢 的 项目!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Pick the one you love most!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
projects=[
    ("🌀","太阳系 转盘","Solar System Wheel","8 行星 围着 太阳 转",SUN),
    ("📏","行星 大小 比较","Planet Size Compare","用 不同 大小 的 圆 排队",NEBULA),
    ("✉️","我 的 太空 明信片","My Space Postcard","写 / 画 — 寄 给 外星人!",STAR),
]
for i,(em,cn,en,d,cl) in enumerate(projects):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(3.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.20,em,sz=80,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.85,0.42,cn,sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.38,2.85,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.85,d,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.90,9.2,0.30,"📕 完成 项目 后 — 一起 做 Day 1 Booklet!",sz=12,b=True,c=MARS,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3 分钟 — 介绍 项目:\n• 4 队 — 每 队 选 一 个 项目 (or 全班 一起 做 同一个)\n• 材料: 纸盘 / 彩笔 / 圆形 不同 大小 (橡皮、硬币、纸杯)\n• 老师 准备: 各 行星 颜色 + 大小 提示")

# ============================================================
# 18. PROJECT DETAILS — 太阳系转盘
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🌀 项目 1 · 太阳系 转盘  Solar System Wheel",SUN)
ib(s,0.4,0.90,4.5,3.95,"🖼️ 转盘 示例图\nSolar wheel sample")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(3.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=SUN; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=SUN; head.line.fill.background()
tb(s,5.25,0.97,4.30,0.40,"📝 怎么 做  How to Make",sz=14,b=True,c=WHITE)
steps=[
    ("1️⃣","拿 一个 大 纸盘 — 中间 画 太阳"),
    ("2️⃣","在 太阳 周围 画 8 个 圈 (行星 轨道)"),
    ("3️⃣","在 每个 圈 上 画 一个 行星 + 写 名字"),
    ("4️⃣","用 不同 颜色 — 让 每个 行星 都 漂亮"),
    ("5️⃣","写 上 你 的 名字 — 这 是 你 的 太空!"),
]
for i,(num,txt) in enumerate(steps):
    y=1.55+i*0.62
    tb(s,5.25,y,0.40,0.40,num,sz=16,b=True,c=SUN)
    tb(s,5.75,y+0.04,3.80,0.35,txt,sz=11,c=DARK)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.95),Inches(9.2),Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb=SUN; tip.line.fill.background()
tb(s,0.55,5.00,9.0,0.30,"💡 提示: 不会画 — 可以 用 贴纸 / 剪 圆形 纸",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.30,9.0,0.18,"Tip: stickers + paper cutouts are OK!",sz=8,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"项目 1 · 25 分钟:\n• 材料: 纸盘 (paper plate) / 彩笔 / 贴纸 / 剪刀\n• 强调: 顺序 是 — 水金地火 木土天海!\n• 老师 走动 — 帮 写 行星 名字 (尤其 K-1)")

# ============================================================
# 19. PROJECT DETAILS — 行星大小比较 + 太空明信片 (combined quick)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📏✉️ 项目 2 + 3  Two More Options",NEBULA)
# Left: planet size compare
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(4.10))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=NEBULA; left.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,0.55,0.97,4.30,0.40,"📏 项目 2 · 行星 大小 比较",sz=14,b=True,c=WHITE)
tb(s,0.55,1.55,4.30,0.32,"用 不同 大小 的 圆 排成 一排 —",sz=12,b=True,c=DARK)
tb(s,0.55,1.90,4.30,0.32,"看 谁 大 谁 小!",sz=12,b=True,c=DARK)
tb(s,0.55,2.35,4.30,0.30,"🌟 材料: 不同 大小 的 圆 (橡皮、纸杯 底、硬币)",sz=10,c=NEBULA)
tb(s,0.55,2.85,4.30,0.32,"🥇 木星 最大 (像 大 西瓜!)",sz=11,b=True,c=DARK)
tb(s,0.55,3.20,4.30,0.32,"🥈 土星 第二 (像 大 苹果)",sz=11,b=True,c=DARK)
tb(s,0.55,3.55,4.30,0.32,"🌍 地球 像 小 葡萄",sz=11,b=True,c=DARK)
tb(s,0.55,3.90,4.30,0.32,"☿ 水星 最小 (像 豌豆!)",sz=11,b=True,c=DARK)
tb(s,0.55,4.40,4.30,0.30,"在 大白纸 上 排 — 每个 圆 写 名字",sz=10,c=NEBULA)
# Right: space postcard
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(4.10))
right.fill.solid(); right.fill.fore_color.rgb=WHITE; right.line.color.rgb=STAR; right.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=STAR; head.line.fill.background()
tb(s,5.20,0.97,4.30,0.40,"✉️ 项目 3 · 我 的 太空 明信片",sz=14,b=True,c=DARK)
tb(s,5.20,1.55,4.30,0.32,"假装 你 在 一个 行星 上 —",sz=12,b=True,c=DARK)
tb(s,5.20,1.90,4.30,0.32,"给 地球 上 的 家人 写 明信片!",sz=12,b=True,c=DARK)
tb(s,5.20,2.35,4.30,0.30,"🌟 材料: 卡纸 / 彩笔 / 贴纸",sz=10,c=STAR)
tb(s,5.20,2.85,4.30,0.32,"📝 句型:",sz=11,b=True,c=DARK)
tb(s,5.20,3.20,4.30,0.32,"「我 在 ___ 行星 — 这里 ___」",sz=12,b=True,c=STAR)
tb(s,5.20,3.55,4.30,0.32,"「我 看到 ___, 我 想 ___」",sz=12,b=True,c=STAR)
tb(s,5.20,3.95,4.30,0.30,"画 你 看到 的 — 太阳? 月亮? 外星人?",sz=10,c=GRAY)
tb(s,5.20,4.40,4.30,0.30,"反 面 画 邮票 + 写 收件人",sz=10,c=STAR)
n+=1; pn(s,n)
notes(s,"项目 2 + 3 · 25 分钟:\n• 让 学生 选 — 哪 个 项目 更 想 做\n• 项目 3 (明信片) 适合 喜欢 写字 的 学生\n• 老师 帮 K-1 学生 写 (老师 写 学生 画)")

# ============================================================
# 20. BOOKLET — Day 1
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📕 Day 1 Booklet  My Solar System Booklet",EARTH)
tb(s,0.4,0.85,9.2,0.32,"把 今天 学到 的 都 画 / 写 进 你 的 小书!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Put everything you learned into your own little book!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
pages=[
    ("1️⃣","封面","Cover","名字 + 太阳 + 星星",SUN),
    ("2️⃣","太阳系","Solar System","画 9 个 球",NEBULA),
    ("3️⃣","最 喜欢","Favorite","你 的 行星 + 为什么",STAR),
    ("4️⃣","太空 词","Space Words","太阳 / 地球 / 宇宙",EARTH),
]
for i,(num,cn,en,d,cl) in enumerate(pages):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.22),Inches(3.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.65,2.12,0.45,num,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.10,2.12,0.42,cn,sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.50,2.12,0.26,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.20,2.85,1.82,1.30,"🖍️ 画 这里")
    tb(s,x+0.05,4.22,2.12,0.36,d,sz=10,b=True,c=cl,a=PP_ALIGN.CENTER)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.80),Inches(9.2),Inches(0.65))
tip.fill.solid(); tip.fill.fore_color.rgb=EARTH; tip.line.fill.background()
tb(s,0.55,4.85,9.0,0.30,"📕 老师 带 着 一起 做 — 一页 一页 翻!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.15,9.0,0.22,"Teacher walks through page by page!",sz=9,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"15-20 分钟 — Day 1 Booklet:\n• 老师 把 booklet 预先 printed (4 页 折一下 就行)\n• 老师 投影 — 自己 也 做 一本 — 学生 跟着 做\n• 一页 一页 — 不 急\n• 完成 后 — 收 起来 (Day 5 展示 用)")

# ============================================================
# 21. SHARE & CLOSE
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎤 分享 + 再见!  Share + Goodbye!",COSMIC)
tb(s,0.4,0.85,9.2,0.32,"今天 你 学 到 了 什么 最 好玩 的?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"What was the most fun thing you learned today?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(9.2),Inches(1.65))
sh.fill.solid(); sh.fill.fore_color.rgb=NIGHT; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,0.55,1.80,9.0,0.40,"💬 句型:",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.25,9.0,0.45,"「我 最 喜欢 ___」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.75,9.0,0.45,"「我 想 去 ___ 行星」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
preview=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.50),Inches(9.2),Inches(1.80))
preview.fill.solid(); preview.fill.fore_color.rgb=WHITE; preview.line.color.rgb=COSMIC; preview.line.width=Pt(2.5)
tb(s,0.55,3.60,9.0,0.40,"🔮 下次 见 (Day 2):",sz=14,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.55,4.05,9.0,0.40,"⭐ 星座 + 银河 + 牛郎织女 的 故事!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.55,9.0,0.30,"Constellations + Milky Way + the Cowherd & Weaver story",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,4.92,9.0,0.30,"👋 晚上 抬头 看 看 天空 — 你 能 找到 月亮 吗?",sz=11,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5 分钟 — 分享 + 收尾:\n• 3-4 个 学生 用 句型 分享\n• 提示 Day 2: 星座 + 中国 神话\n• 「家庭 作业」: 今晚 抬头 看 一下 天空")

# === SAVE ===
import os
out=os.path.join(os.path.dirname(__file__),"day1_solar_system.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
