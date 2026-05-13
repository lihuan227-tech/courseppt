#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
仰望星空 · Day 4: 外星人 是否 存在  Do Aliens Exist?
Story anchor: 来自 火星 的 一 封 信 (teacher-written)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random, os

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

NIGHT  = RGBColor(0x0D,0x1B,0x3E)
COSMIC = RGBColor(0x6A,0x1B,0x9A)
STAR   = RGBColor(0xF5,0xC2,0x42)
GOLD   = RGBColor(0xFF,0xB7,0x00)
EARTH  = RGBColor(0x1E,0x88,0xE5)
RED    = RGBColor(0xC8,0x25,0x3E)
MARS   = RGBColor(0xD8,0x43,0x15)
PINK   = RGBColor(0xEC,0x40,0x7A)
NEBULA = RGBColor(0x7B,0x1F,0xA2)
ALIEN  = RGBColor(0x66,0xBB,0x6A)
SKY    = RGBColor(0x42,0xA5,0xF5)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=IMGBG; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=STAR,a=PP_ALIGN.CENTER)
    for x,y in [(0.8,4.7),(1.8,4.5),(7.8,4.5),(8.6,4.7)]:
        d=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(x),Inches(y),Inches(0.35),Inches(0.35))
        d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
    return s
def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mx=x+size/2; my=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mx-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(my-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

n=0

# Cover
s=ns(); bg(s,NIGHT)
random.seed(23)
for _ in range(40):
    x=random.uniform(0.3,9.7); y=random.uniform(0.3,5.3); sz=random.uniform(0.06,0.16)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
tb(s,0.5,0.4,9,0.6,"🌌 仰望星空  Looking Up at the Stars",sz=28,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.0,9,0.45,"Day 4 · 外星人 是否 存在?  Do Aliens Exist?",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# UFO + Alien + Earth
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(4.0),Inches(2.0),Inches(2.0),Inches(2.0))
sh.fill.solid(); sh.fill.fore_color.rgb=ALIEN; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,4.0,2.50,2.0,0.6,"👽",sz=70,a=PP_ALIGN.CENTER)
tb(s,4.0,3.20,2.0,0.4,"外星人?",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.0,2.3,2.5,1.5,"🛸",sz=80,a=PP_ALIGN.CENTER)
tb(s,7.0,2.3,2.5,1.5,"🌍",sz=80,a=PP_ALIGN.CENTER)
tb(s,0.5,4.55,9,0.40,"✉️ 来自 火星 的 一 封 信",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,4.95,9,0.25,"A Letter from Mars (teacher-written)",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 4 开场:\n• 「今天 — 我们 来 想 一个 大 问题: 外星人 真的 存在 吗?」\n• 老师 神秘 地 拿出 一 封 信 ……\n• 「这 是 — 来自 火星 的 信!」")

# Session 1 Divider
s=div("Session 1  上午 11:00–11:45","👽 故事课 · 来自 火星 的 一 封 信  ·  45 min",ALIEN,"🛸"); n+=1; pn(s,n)

# Learning Goals
s=ns(); bg(s,CREAM); hb(s,"🎯 今天的学习目标  Today's Learning Goals",NIGHT)
tb(s,0.4,0.85,9.2,0.30,"上完这节课, 你会……",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
goals=[
    ("1","✉️","通过 老师 自编 故事 — 想象 外星 生命",ALIEN),
    ("2","🧠","学会 区分 「事实」 + 「想象」",STAR),
    ("3","🙋","提问 + 表达 你 自己 的 观点",NEBULA),
    ("4","🔭","知道 宇宙 还有 很多 未知 等 我们 探索",COSMIC),
]
for i,(num,em,text,cl) in enumerate(goals):
    y=1.30+i*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.18),Inches(0.50),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.22,0.50,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.18,0.60,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.90,y+0.20,7.6,0.55,text,sz=14,b=True,c=DARK)
n+=1; pn(s,n)

# Hook — Do aliens exist?
s=ns(); bg(s,CREAM); hb(s,"🤔 你 觉得 外星人 存在 吗?  Do You Think Aliens Exist?",ALIEN)
tb(s,0.4,0.85,9.2,0.32,"听听 你 的 想法 — 没有 错 答案!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"What do YOU think? No wrong answers!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
opts=[
    ("✅","存在!","Yes!","宇宙 太 大 — 肯定 还 有 别的 生命",ALIEN),
    ("❌","不 存在","No","只有 地球 有 生命",EARTH),
    ("🤷","不知道","Don't know","还没 找到 — 但 也许 有",COSMIC),
]
for i,(em,cn,en,d,cl) in enumerate(opts):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(3.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.20,em,sz=78,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.85,0.40,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.40,2.85,0.28,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.85,d,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.80),Inches(9.2),Inches(0.65))
prompt.fill.solid(); prompt.fill.fore_color.rgb=NIGHT; prompt.line.color.rgb=STAR; prompt.line.width=Pt(2.5)
tb(s,0.55,4.85,9.0,0.30,"🗳️ 举手 投票! 然后 说说 — 为什么?",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,5.18,9.0,0.20,"Vote with hands — and tell us why!",sz=8,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"4 分钟 — Hook:\n• 投票 — 让 学生 选 一边\n• 关键: 没有 错 答案!\n• 问 2-3 个 学生: 「为什么 你 这么 想?」")

# Letter from Mars (Story)
s=ns(); bg(s,CREAM); hb(s,"✉️ 来自 火星 的 一 封 信  Letter from Mars",MARS)
# letter card
letter=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(0.95),Inches(9.0),Inches(4.50))
letter.fill.solid(); letter.fill.fore_color.rgb=WARM; letter.line.color.rgb=MARS; letter.line.width=Pt(3)
tb(s,0.7,1.05,8.6,0.42,"📬 亲爱 的 地球 朋友:",sz=16,b=True,c=MARS)
tb(s,0.7,1.50,8.6,0.30,"Dear Earth Friends,",sz=10,c=GRAY)
lines=[
    "🛸 我 是 火星人 Zorp — 今年 100 岁 (在 火星 算 小朋友)。",
    "🌍 我 们 一直 在 看 你们 — 你们 的 蓝色 星球 太 漂亮 了!",
    "🤔 我 有 一个 问题: 你们 也 在 找 我们 吗?",
    "🌌 宇宙 有 几 千 亿 个 星系 — 肯定 还 有 别的 生命!",
    "📡 也许 有 一 天 — 我们 会 互相 找到。",
    "💌 在 那 之前 — 请 保护 好 你们 的 地球。",
    "✨ 你的 火星 朋友, Zorp",
]
for i,l in enumerate(lines):
    tb(s,0.7,1.90+i*0.43,8.6,0.40,l,sz=12,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"5-7 分钟 — 读 火星 来信:\n• 老师 用 神秘 的 语气 读\n• 中间 停 一下 — 让 学生 反应\n• 提问: 「Zorp 觉得 外星人 存在 吗?」 → 存在!\n• 引导: 「我们 还 没 找到 — 但 也许 有 一 天 ……」")

# ============================================================
# Letter — pre-reading predictions
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔮 拆 信 之前 — 想 一 想!  Before We Read",ALIEN)
tb(s,0.4,0.85,9.2,0.34,"信 封 写 着 「来自 火星」 — 你 猜 信 里 说 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Envelope says 'From Mars' — what's in the letter?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
preds=[
    ("👽","写 信 的 人 长 什么 样?","What do they look like?",ALIEN),
    ("💬","他们 说 什么 语言?","What language do they speak?",NEBULA),
    ("❓","他们 想 跟 我们 说 什么?","What do they want to tell us?",MARS),
]
for i,(em,q,en,cl) in enumerate(preds):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.0,em,sz=66,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,2.85,0.50,q,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.30,2.85,0.40,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.80,2.75,0.55,"💡 大胆 猜!",sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
tps=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.55),Inches(9.2),Inches(0.95))
tps.fill.solid(); tps.fill.fore_color.rgb=NIGHT; tps.line.color.rgb=STAR; tps.line.width=Pt(2.5)
tb(s,0.55,4.62,9.0,0.30,"👥 Think-Pair-Share: 跟 同桌 猜 (1 分钟)",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.26,"Turn to a partner — guess what's inside!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.24,"💬 「信 里 也许 说 ___」",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🔮 PRE-READING · 3 分钟:\n• 戏剧化: 「老师 收到 一封 信 — 来自 火星!」\n• 让 学生 猜 — 不评判\n• 收集 想法 — 然后 「我们 来 拆 信!」")

# ============================================================
# Listening missions during letter
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"👀 拆 信 时 — 你 的 任务!  While You Listen",STAR)
tb(s,0.4,0.85,9.2,0.34,"3 个 任务 — 一边 听 一边 留心!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"3 missions — listen + observe!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
obs=[
    ("1","🛸","Zorp 是 谁? 多 大 年纪?","Who is Zorp? How old?",ALIEN),
    ("2","🌍","他 对 地球 说 了 什么?","What did he say about Earth?",EARTH),
    ("3","💌","他 提了 几 个 问题?","How many questions did he ask?",NEBULA),
]
for i,(num,em,cn,en,cl) in enumerate(obs):
    y=1.55+i*1.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.95))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.20),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.27,0.55,0.40,num,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.18,0.80,0.55,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,2.30,y+0.14,7.0,0.40,cn,sz=15,b=True,c=cl)
    tb(s,2.30,y+0.54,7.0,0.30,en,sz=10,c=GRAY)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.60))
tip.fill.solid(); tip.fill.fore_color.rgb=STAR; tip.line.fill.background()
tb(s,0.55,4.93,9.0,0.30,"👂 用 心 听! 念 完 — 我们 一起 讨论 + 回信",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,5.23,9.0,0.22,"Listen well — we'll discuss + reply after.",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"👀 DURING-READING · 1 分钟 setup + 5 分钟 念信:\n• 念 3 个 任务\n• 老师 念 信 (神秘 语气)\n• 念 完 → 下页 讨论")

# ============================================================
# Post-letter discussion (6 Qs)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"💭 拆 信 后 — 一起 讨论!  After Reading",ALIEN)
tb(s,0.4,0.85,9.2,0.30,"选 1-2 个 问题 — 全班 / Think-Pair-Share",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.15,9.2,0.22,"Pick 1-2 — class / Think-Pair-Share",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("👽","你 觉得 Zorp 长 什么 样?","What does Zorp look like?"),
    ("💌","他 说 「保护 地球」 — 为什么?","Why did he say 'protect Earth'?"),
    ("🤝","如果 你 见到 Zorp — 你 说 什么?","If you met Zorp — what would you say?"),
    ("🌍","信 是 真的 还是 假的? 怎么 知道?","Real or fake letter? How do we know?"),
    ("🛸","为什么 我们 还 没 找到 外星人?","Why haven't we found aliens yet?"),
    ("💭","你 想 给 Zorp 回 什么 信?","What letter would YOU write back?"),
]
for i,(em,cn,en) in enumerate(qs):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.45+row*1.80
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.65))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=ALIEN; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.42),Inches(0.42))
    nb.fill.solid(); nb.fill.fore_color.rgb=ALIEN; nb.line.fill.background()
    tb(s,x+0.10,y+0.14,0.42,0.34,str(i+1),sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.60,y+0.05,0.55,0.50,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,y+0.60,2.70,0.65,cn,sz=11,b=True,c=DARK)
    tb(s,x+0.15,y+1.25,2.70,0.32,en,sz=8,c=GRAY)
n+=1; pn(s,n)
notes(s,"💭 POST-LETTER · 5-7 分钟:\n• 选 2-3 题\n• Q3 + Q6 启发 想象\n• Q4 — 科学 思维 (real vs imagination)\n• 鼓励 各种 答案 — 没 标准")

# Fact vs Imagination
s=ns(); bg(s,CREAM); hb(s,"🧠 事实 vs 想象  Fact vs Imagination",NEBULA)
tb(s,0.4,0.85,9.2,0.32,"哪些 是 真的? 哪些 是 想 出来 的?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"What's real? What's imagined?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Two columns
fact=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(4.55),Inches(3.20))
fact.fill.solid(); fact.fill.fore_color.rgb=WHITE; fact.line.color.rgb=EARTH; fact.line.width=Pt(3)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(4.55),Inches(0.55))
head.fill.solid(); head.fill.fore_color.rgb=EARTH; head.line.fill.background()
tb(s,0.55,1.62,4.30,0.40,"✅ 事实 Fact (真 的)",sz=14,b=True,c=WHITE)
facts=[
    "🌌 宇宙 有 几 千 亿 个 星系",
    "🔴 火星 有 水 (冰)",
    "🛰️ 人类 发射 过 探测器 找 生命",
    "🌍 地球 是 唯一 已知 有 生命 的 星球",
]
for i,f in enumerate(facts):
    tb(s,0.55,2.25+i*0.55,4.30,0.50,f,sz=11,b=True,c=DARK)
imag=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.55),Inches(4.55),Inches(3.20))
imag.fill.solid(); imag.fill.fore_color.rgb=WHITE; imag.line.color.rgb=NEBULA; imag.line.width=Pt(3)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.55),Inches(4.55),Inches(0.55))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,5.20,1.62,4.30,0.40,"💭 想象 Imagination",sz=14,b=True,c=WHITE)
imags=[
    "👽 外星人 长 什么 样?",
    "🛸 飞碟 真的 来过 地球?",
    "🌟 别的 星球 有 高 智能 文明?",
    "📡 他们 也 想 找 我们 吗?",
]
for i,f in enumerate(imags):
    tb(s,5.20,2.25+i*0.55,4.30,0.50,f,sz=11,b=True,c=DARK)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.90),Inches(9.2),Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb=GOLD; tip.line.fill.background()
tb(s,0.55,4.97,9.0,0.32,"💡 两个 都 重要 — 事实 帮我们 学习, 想象 帮我们 创造!",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"4 分钟 — 事实 vs 想象:\n• 一对 一对 念\n• 关键 message: 「都 重要! 事实 = 科学, 想象 = 故事 + 发明」\n• 互动: 让 学生 再 加 一个 事实 + 一个 想象")

# Session 1 wrap
s=ns(); bg(s,CREAM); hb(s,"🎤 你 怎么 想?  What Do You Think?",ALIEN)
qs=[
    ("👽","如果 外星人 来 地球 — 你 想 问 他们 什么?",ALIEN),
    ("🤔","你 觉得 他们 长 什么 样? 为什么?",NEBULA),
    ("🌍","如果 你 写 回信 给 Zorp — 你 想 说 什么?",MARS),
]
tb(s,0.4,0.85,9.2,0.32,"自由 提问 + 表达 你 自己 的 观点!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,q,cl) in enumerate(qs):
    y=1.40+i*1.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.00))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,0.55,y+0.22,0.80,0.60,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.30,8.0,0.50,q,sz=14,b=True,c=cl)
n+=1; pn(s,n)

# Session 2 Divider
s=div("Session 2  下午 2:00–2:45","📚 词汇课 · 我会认 + 我会写  ·  45 min",EARTH,"📖"); n+=1; pn(s,n)

# Review
s=ns(); bg(s,CREAM); hb(s,"🔁 早上 学了 什么?",EARTH)
items=[
    ("✉️","火星 来信","Letter from Mars","Zorp 的 故事",MARS),
    ("👽","外星人?","Aliens?","存在 / 不存在 / 不知道",ALIEN),
    ("🧠","事实 vs 想象","Fact vs Imagine","两个 都 重要!",NEBULA),
    ("🌌","宇宙 未知","Unknown universe","还 有 好多 等 我们 发现",COSMIC),
]
tb(s,0.4,0.85,9.2,0.32,"想 一 想 — 早上 我们 学了 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,cn,en,d,cl) in enumerate(items):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.40+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.20,1.00,1.20,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.18,3.25,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.20,y+0.58,3.25,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.92,3.25,0.70,d,sz=11,b=True,c=DARK)
n+=1; pn(s,n)

# 我会认 5 words
s=ns(); bg(s,CREAM); hb(s,"📖 我会认  I Can Read",STAR)
tb(s,0.4,0.85,9.2,0.32,"5 个 词 — 一起 读!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
words=[
    ("👽","外星人","wài xīng rén","alien",ALIEN),
    ("🌱","生命","shēng mìng","life",GREEN_OK := RGBColor(0x38,0x8E,0x3C)),
    ("📡","信号","xìn hào","signal",SKY),
    ("🔍","发现","fā xiàn","discover",STAR),
    ("💭","猜想","cāi xiǎng","guess",NEBULA),
]
for i,(em,cn,py,en,cl) in enumerate(words):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.78),Inches(3.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,1.70,0.90,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,1.70,0.55,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.25,1.70,0.35,py,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.62,1.70,0.30,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.05,1.60,0.85,"跟读\n3 遍",sz=11,b=True,c=cl,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

# 我会写 — 外星人, 生命
def write_slide(emoji,word_cn,word_en,chars,color):
    s=ns(); bg(s,CREAM); hb(s,f"✏️ 我会写 · {word_cn}  I Can Write · {word_en}",color)
    tb(s,0.4,0.85,9.2,0.36,f"{emoji} 一起来写「{word_cn}」!",sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.25,9.2,0.26,f"Practice writing {word_cn} ({word_en})",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    if len(chars)==3:
        for i,(ch,py,_,_) in enumerate(chars):
            tianzi(s,0.55+i*1.65,1.65,1.55,ch,color,pinyin=py,char_sz=90)
    elif len(chars)==2:
        tianzi(s,0.55,1.65,2.20,chars[0][0],color,pinyin=chars[0][1],char_sz=120)
        tianzi(s,2.95,1.65,2.20,chars[1][0],color,pinyin=chars[1][1],char_sz=120)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.45,1.72,4.10,0.40,"✏️ 怎么写  How to Write",sz=13,b=True,c=WHITE)
    spacing = 0.62 if len(chars)==3 else 0.95
    for i,(ch,py,hint_cn,hint_en) in enumerate(chars):
        y=2.30+i*spacing
        tb(s,5.45,y,4.10,0.30,f"📐「{ch}」 — {py}",sz=12,b=True,c=DARK)
        tb(s,5.45,y+0.28,4.10,0.26,hint_cn,sz=9,b=True,c=color)
    pf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
    pf.fill.solid(); pf.fill.fore_color.rgb=WARM; pf.line.color.rgb=color; pf.line.width=Pt(2)
    tb(s,0.55,4.72,9.0,0.32,f"📝 在 田字格 里 写 3 遍",sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.55,5.08,9.0,0.32,f"💬 「我 会 写「{word_cn}」!」",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    return s

s=write_slide("👽","外星人","Alien",[
    ("外","wài","5 笔 — 「夕」+「卜」","Evening + divination"),
    ("星","xīng","9 笔 — 「日」+「生」","Sun + life"),
    ("人","rén","2 笔 — 像 一 个 人 站 着","Looks like a standing person"),
],ALIEN); n+=1; pn(s,n)

s=write_slide("🌱","生命","Life",[
    ("生","shēng","5 笔 — 像 一 颗 小草 生长","Sprouting plant"),
    ("命","mìng","8 笔 — 「人」+「一」+「叩」","Person + line + 叩"),
],GREEN_OK); n+=1; pn(s,n)

# Session 3 Divider
s=div("Session 3  下午 3:00–4:30","🛠️ 项目 · 外星人 创意 工作坊  ·  90 min",ALIEN,"👽"); n+=1; pn(s,n)

# Projects overview
s=ns(); bg(s,CREAM); hb(s,"🛠️ 3 个 项目  3 Projects",ALIEN)
tb(s,0.4,0.85,9.2,0.32,"展开 你 的 想象 — 选 一个!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
projects=[
    ("👽","设计 外星人","Design Alien","几 只 眼? 什么 颜色?",ALIEN),
    ("🛸","外星 飞船","Alien Spaceship","你 的 UFO 是 什么 样?",NEBULA),
    ("📡","发 消息 给 外星人","Send Message","画 + 写 — 让 他们 收 到!",STAR),
]
for i,(em,cn,en,d,cl) in enumerate(projects):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(3.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.20,em,sz=78,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.85,0.42,cn,sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.38,2.85,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.85,d,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

# Project 1 details
s=ns(); bg(s,CREAM); hb(s,"👽 项目 1 · 设计 你 的 外星人  Design YOUR Alien",ALIEN)
ib(s,0.4,0.90,4.5,3.95,"🖼️ 外星人 创意 示例")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(3.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=ALIEN; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=ALIEN; head.line.fill.background()
tb(s,5.25,0.97,4.30,0.40,"📝 想 一 想",sz=14,b=True,c=WHITE)
qs=[
    "1️⃣ 几 只 眼睛?",
    "2️⃣ 什么 颜色?",
    "3️⃣ 几 条 腿 / 触角?",
    "4️⃣ 它 吃 什么?",
    "5️⃣ 它 住 在 哪个 星球?",
    "6️⃣ 给 它 起 一个 名字!",
]
for i,q in enumerate(qs):
    y=1.55+i*0.50
    tb(s,5.25,y,4.30,0.40,q,sz=12,b=True,c=DARK)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.95),Inches(9.2),Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb=ALIEN; tip.line.fill.background()
tb(s,0.55,5.02,9.0,0.32,"💡 没有 错 答案 — 越 奇怪 越 好玩!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"项目 1 · 25 分钟:\n• 材料: 彩纸 / 彩笔 / 贴纸 / 棉球\n• 老师 示范: 画 一个 紫色 3 眼 外星人\n• 鼓励 学生 自由 创作 — 不 评 像 不 像")

# Project 2 + 3
s=ns(); bg(s,CREAM); hb(s,"🛸📡 项目 2 + 3",NEBULA)
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(4.10))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=NEBULA; left.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,0.55,0.97,4.30,0.40,"🛸 项目 2 · 外星 飞船 设计",sz=14,b=True,c=WHITE)
items=[
    "🌟 你 的 UFO 是 什么 样?",
    "1️⃣ 画 形状 (圆? 三角? 像 鸡蛋?)",
    "2️⃣ 加 灯 + 窗户 + 装饰",
    "3️⃣ 几 个 外星 乘客?",
    "4️⃣ 它 能 飞 多 快? 烧 什么 燃料?",
    "5️⃣ 给 它 起 一个 名字!",
]
for i,line in enumerate(items):
    tb(s,0.55,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(4.10))
right.fill.solid(); right.fill.fore_color.rgb=WHITE; right.line.color.rgb=STAR; right.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=STAR; head.line.fill.background()
tb(s,5.20,0.97,4.30,0.40,"📡 项目 3 · 给 外星人 发 消息",sz=14,b=True,c=DARK)
parts=[
    "🌟 假装 你 是 NASA — 发 一条 消息 给 Zorp!",
    "1️⃣ 画 地球 + 你 自己 + 妈妈 / 朋友",
    "2️⃣ 用 简单 图 解释 你 是 谁",
    "3️⃣ 写 一句话 中文 + 一句话 英文",
    "「你 好! 我 是 ___」",
    "「请 来 我们 地球 玩!」",
]
for i,line in enumerate(parts):
    tb(s,5.20,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"30 分钟:\n• 提示: NASA 真的 发过 Voyager 金 盘 给 外星人!\n• 让 学生 觉得 — 「这 是 真 的 工作!」")

# Share & close
s=ns(); bg(s,CREAM); hb(s,"🎤 分享 + 再见!",COSMIC)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(2.10))
sh.fill.solid(); sh.fill.fore_color.rgb=NIGHT; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,0.55,1.05,9.0,0.40,"💬 句型:",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,1.55,9.0,0.50,"「我 的 外星人 叫 ___」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.10,9.0,0.50,"「它 来自 ___ 星球, 它 会 ___」",sz=20,b=True,c=STAR,a=PP_ALIGN.CENTER)
preview=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.25),Inches(9.2),Inches(2.05))
preview.fill.solid(); preview.fill.fore_color.rgb=WHITE; preview.line.color.rgb=COSMIC; preview.line.width=Pt(2.5)
tb(s,0.55,3.40,9.0,0.40,"🔮 下次 见 (Day 5):",sz=14,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.55,3.85,9.0,0.40,"🌟 我 的 太空 梦想 + Final Showcase!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.30,9.0,0.30,"My Space Dream + Final Showcase",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,4.70,9.0,0.30,"👋 想 一想: 你 的 太空 梦想 是 什么?",sz=11,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

out=os.path.join(os.path.dirname(__file__),"day4_aliens.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
