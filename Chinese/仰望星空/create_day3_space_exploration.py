#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
仰望星空 · Day 3: 外太空探索 — 登月 和 火星 计划  Space Exploration
Picture book anchor: 如果 你 决定 去 月球  If You Decide to Go to the Moon
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random, os

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

NIGHT  = RGBColor(0x0D,0x1B,0x3E)
COSMIC = RGBColor(0x6A,0x1B,0x9A)
STAR   = RGBColor(0xF5,0xC2,0x42)
GOLD   = RGBColor(0xFF,0xB7,0x00)
EARTH  = RGBColor(0x1E,0x88,0xE5)
RED    = RGBColor(0xC8,0x25,0x3E)
MARS   = RGBColor(0xD8,0x43,0x15)
NEBULA = RGBColor(0x7B,0x1F,0xA2)
SKY    = RGBColor(0x42,0xA5,0xF5)
MOON_C = RGBColor(0xB0,0xBE,0xC5)
GREEN  = RGBColor(0x2E,0x7D,0x32)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=NIGHT,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=IMGBG; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=STAR,a=PP_ALIGN.CENTER)
    for x,y in [(0.8,4.7),(1.8,4.5),(7.8,4.5),(8.6,4.7)]:
        d=s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(x),Inches(y),Inches(0.35),Inches(0.35))
        d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
    return s
def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mx=x+size/2; my=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mx-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(my-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

n=0

# Cover
s=ns(); bg(s,NIGHT)
random.seed(17)
for _ in range(40):
    x=random.uniform(0.3,9.7); y=random.uniform(0.3,5.3); sz=random.uniform(0.06,0.16)
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(sz),Inches(sz))
    d.fill.solid(); d.fill.fore_color.rgb=STAR; d.line.fill.background()
tb(s,0.5,0.4,9,0.6,"🌌 仰望星空  Looking Up at the Stars",sz=28,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,1.0,9,0.45,"Day 3 · 太空 探索  Space Exploration",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# Earth, Moon, Mars
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(1.5),Inches(2.0),Inches(1.7),Inches(1.7))
sh.fill.solid(); sh.fill.fore_color.rgb=EARTH; sh.line.color.rgb=GREEN; sh.line.width=Pt(3)
tb(s,1.5,2.40,1.7,0.6,"🌍",sz=44,a=PP_ALIGN.CENTER)
tb(s,1.5,3.55,1.7,0.4,"地球",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(4.15),Inches(2.0),Inches(1.7),Inches(1.7))
sh.fill.solid(); sh.fill.fore_color.rgb=MOON_C; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,4.15,2.40,1.7,0.6,"🌕",sz=44,a=PP_ALIGN.CENTER)
tb(s,4.15,3.55,1.7,0.4,"月球",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(6.8),Inches(2.0),Inches(1.7),Inches(1.7))
sh.fill.solid(); sh.fill.fore_color.rgb=MARS; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,6.8,2.40,1.7,0.6,"🔴",sz=44,a=PP_ALIGN.CENTER)
tb(s,6.8,3.55,1.7,0.4,"火星",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
# rocket
tb(s,8.7,2.10,1.2,1.6,"🚀",sz=70,a=PP_ALIGN.CENTER)
tb(s,0.5,4.55,9,0.40,"📖 绘本: 如果 你 决定 去 月球",sz=15,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.5,4.95,9,0.25,"Picture book: If You Decide to Go to the Moon",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 3 开场:\n• 「今天 — 我们 真的 要 去 太空 旅行!」\n• 介绍 绘本: 如果 你 决定 去 月球\n• 准备 道具: 火箭 模型 / 宇航员 图片")

# Session 1 divider
s=div("Session 1  上午 11:00–11:45","🚀 故事课 · 去 月球 + 火星 大冒险  ·  45 min",NEBULA,"🌕"); n+=1; pn(s,n)

# Learning Goals
s=ns(); bg(s,CREAM); hb(s,"🎯 今天的学习目标  Today's Learning Goals",NIGHT)
tb(s,0.4,0.85,9.2,0.30,"上完这节课, 你会……  By the end, you'll be able to…",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
goals=[
    ("1","🌕","通过 绘本 — 想象 真实 太空 旅行 会 遇到 什么 问题",MOON_C),
    ("2","🚀","认识 宇航员、火箭、月球、火星",MARS),
    ("3","💡","知道 人类 为什么 要 探索 太空",STAR),
    ("4","🌟","对 未来 科技 + 探索 充满 兴趣!",NEBULA),
]
for i,(num,em,text,cl) in enumerate(goals):
    y=1.30+i*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.18),Inches(0.50),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.22,0.50,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.18,0.60,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.90,y+0.20,7.6,0.55,text,sz=14,b=True,c=DARK)
n+=1; pn(s,n)

# Hook — Imagine packing for space
s=ns(); bg(s,CREAM); hb(s,"🎒 准备 上 太空 — 带 什么?  Packing for Space!",MARS)
tb(s,0.4,0.85,9.2,0.32,"假装 你 明天 上 月球 — 你 要 带 哪些 东西?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Imagine you fly to the moon tomorrow — what do you bring?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
items=[
    ("👨‍🚀","宇航服","Spacesuit","没空气, 要 自己 带!",MOON_C),
    ("💧","水 + 食物","Food/water","月球 没有!",EARTH),
    ("🧴","氧气 罐","Oxygen tank","呼吸 用",SKY),
    ("📷","相机","Camera","拍 给 地球 看!",GOLD),
]
for i,(em,cn,en,d,cl) in enumerate(items):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.22),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.12,1.0,em,sz=58,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.85,2.12,0.36,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.22,2.12,0.24,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.55,2.02,0.75,d,sz=10,c=DARK,a=PP_ALIGN.CENTER)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.55),Inches(9.2),Inches(0.95))
prompt.fill.solid(); prompt.fill.fore_color.rgb=NEBULA; prompt.line.color.rgb=STAR; prompt.line.width=Pt(2.5)
tb(s,0.55,4.62,9.0,0.30,"🙋 你 还 想 带 什么? 举手 说说看!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.30,"💬 「我 想 带 ___ 因为 ___ 。」",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,5.25,9.0,0.20,"I want to bring ___ because ___.",sz=8,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"4 分钟 — Hook:\n• 让 学生 想 — 上 太空 要 带 什么?\n• 收 3-5 个 答案\n• 提示: 「月球 没 空气、没 水、没 食物 — 都 要 带!」")

# Picture book intro
s=ns(); bg(s,CREAM); hb(s,"📖 绘本: 如果 你 决定 去 月球",MOON_C)
tb(s,0.4,0.85,9.2,0.40,"作者: Faith McNulty — 真实 又 好玩 的 月球 之 旅",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"A real-yet-playful trip to the moon, by Faith McNulty",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.5,1.65,4.5,3.30,"📚 绘本 封面")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(3.30))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=MOON_C; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(1.65),Inches(4.35),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=MOON_C; head.line.fill.background()
tb(s,5.40,1.72,4.10,0.40,"📚 故事 简介",sz=13,b=True,c=DARK)
parts=[
    "🚀 你 真的 决定 去 月球 — 怎么 准备?",
    "🌑 月球 离 地球 38 万 公里",
    "⏱️ 火箭 飞 3 天 才 到!",
    "🦘 月球 重力 小 — 跳 起来 高 6 倍!",
    "🌡️ 白天 100°C, 晚上 -150°C",
    "🌍 在 月球 看 地球 — 又 蓝 又 亮!",
]
for i,line in enumerate(parts):
    tb(s,5.40,2.25+i*0.42,4.10,0.40,line,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"5-6 分钟 — 讲 月球 之 旅:\n• 念 故事 — 强调 真实 数字 让 学生 「哇!」\n• 互动: 「跳 6 倍 高 — 演 一下!」")

# Why explore space?
s=ns(); bg(s,NIGHT); hb(s,"💡 为 什么 要 探索 太空?  Why Explore?",STAR)
tb(s,0.4,0.85,9.2,0.30,"人类 为什么 想 去 月球 + 火星?",sz=13,b=True,c=STAR,a=PP_ALIGN.CENTER)
reasons=[
    ("🔍","好奇","Curiosity","「那边 是 什么?」",STAR),
    ("🧪","研究","Research","学 更多 宇宙 知识",NEBULA),
    ("🌱","新家","New Home","未来 也许 住 火星?",MARS),
    ("🌍","保护 地球","Save Earth","学 怎么 不 浪费 资源",EARTH),
]
for i,(em,cn,en,d,cl) in enumerate(reasons):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.45+row*1.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.60))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.10,1.00,0.85,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.10,3.25,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.20,y+0.50,3.25,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.85,3.25,0.65,d,sz=12,b=True,c=DARK)
tb(s,0.4,5.10,9.2,0.30,"🌟 1969年 — Apollo 11 第一次 登月!",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"4 分钟:\n• 4 个 原因 — 一个 一个 念\n• 引出: 「1969 年 第一次 登月 — Neil Armstrong」\n• 「现在 — 我们 在 准备 去 火星!」")

# Session 1 wrap
s=ns(); bg(s,CREAM); hb(s,"🎤 一起 想 一 想",NEBULA)
qs=[
    ("🚀","如果 你 是 宇航员 — 你 想 去 月球 还是 火星? 为什么?",MARS),
    ("🌍","在 太空 看 地球 — 你 觉得 它 是 什么 样?",EARTH),
    ("🔮","100 年 后 — 我们 会 住 在 哪里?",NEBULA),
]
tb(s,0.4,0.85,9.2,0.32,"听完 故事 — 你 觉得……",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,q,cl) in enumerate(qs):
    y=1.40+i*1.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.00))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,0.55,y+0.22,0.80,0.60,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.30,8.0,0.50,q,sz=14,b=True,c=cl)
n+=1; pn(s,n)
notes(s,"5 分钟 — 讨论")

# Session 2 Divider
s=div("Session 2  下午 2:00–2:45","📚 词汇课 · 我会认 + 我会写  ·  45 min",EARTH,"📖"); n+=1; pn(s,n)

# Review
s=ns(); bg(s,CREAM); hb(s,"🔁 早上 学了 什么?",EARTH)
items=[
    ("🚀","火箭 + 宇航员","Rockets + astronauts","怎么 去 太空",MARS),
    ("🌕","月球","Moon","跳 6 倍 高!",MOON_C),
    ("🔴","火星","Mars","红色 星球",MARS),
    ("💡","为什么 探索","Why explore","好奇 + 研究 + 新家",STAR),
]
tb(s,0.4,0.85,9.2,0.32,"想 一 想 — 早上 我们 学了 什么?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,cn,en,d,cl) in enumerate(items):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.40+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.20,1.00,1.20,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.18,3.25,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.20,y+0.58,3.25,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.92,3.25,0.70,d,sz=11,b=True,c=DARK)
n+=1; pn(s,n)

# 我会认 5 words
s=ns(); bg(s,CREAM); hb(s,"📖 我会认  I Can Read",STAR)
tb(s,0.4,0.85,9.2,0.32,"5 个 太空 词 — 一起 读!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
words=[
    ("👨‍🚀","宇航员","yǔ háng yuán","astronaut",MOON_C),
    ("🚀","火箭","huǒ jiàn","rocket",MARS),
    ("🌕","月球","yuè qiú","moon",STAR),
    ("🔴","火星","huǒ xīng","Mars",MARS),
    ("🛰️","太空 站","tài kōng zhàn","space station",SKY),
]
for i,(em,cn,py,en,cl) in enumerate(words):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.78),Inches(3.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,1.70,0.90,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,1.70,0.55,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.25,1.70,0.35,py,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.62,1.70,0.30,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.05,1.60,0.85,"跟读\n3 遍",sz=11,b=True,c=cl,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

# 我会写 — 3 words
def write_slide(emoji,word_cn,word_en,chars,color):
    s=ns(); bg(s,CREAM); hb(s,f"✏️ 我会写 · {word_cn}  I Can Write · {word_en}",color)
    tb(s,0.4,0.85,9.2,0.36,f"{emoji} 一起来写「{word_cn}」!",sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.25,9.2,0.26,f"Practice writing {word_cn} ({word_en})",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    if len(chars)==2:
        tianzi(s,0.55,1.65,2.20,chars[0][0],color,pinyin=chars[0][1],char_sz=120)
        tianzi(s,2.95,1.65,2.20,chars[1][0],color,pinyin=chars[1][1],char_sz=120)
    else:
        tianzi(s,1.30,1.65,2.95,chars[0][0],color,pinyin=chars[0][1],char_sz=160)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.45,1.72,4.10,0.40,"✏️ 怎么写  How to Write",sz=13,b=True,c=WHITE)
    for i,(ch,py,hint_cn,hint_en) in enumerate(chars):
        y=2.30+i*0.95
        tb(s,5.45,y,4.10,0.35,f"📐「{ch}」 — {py}",sz=13,b=True,c=DARK)
        tb(s,5.45,y+0.32,4.10,0.30,hint_cn,sz=10,b=True,c=color)
        tb(s,5.45,y+0.58,4.10,0.26,hint_en,sz=9,c=GRAY)
    pf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
    pf.fill.solid(); pf.fill.fore_color.rgb=WARM; pf.line.color.rgb=color; pf.line.width=Pt(2)
    tb(s,0.55,4.72,9.0,0.32,f"📝 在 田字格 里 写 3 遍",sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.55,5.08,9.0,0.32,f"💬 「我 会 写「{word_cn}」!」",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    return s

s=write_slide("🚀","火箭","Rocket",[
    ("火","huǒ","4 笔 — 像 跳动 的 火苗","Like leaping flames"),
    ("箭","jiàn","15 笔 — 竹字头 + 「前」","Bamboo top + 前"),
],MARS); n+=1; pn(s,n)

s=write_slide("🌕","月球","Moon",[
    ("月","yuè","4 笔 — 像 弯弯 的 月亮","Curved like the moon"),
    ("球","qiú","11 笔 — 「王」+「求」","King + 求"),
],MOON_C); n+=1; pn(s,n)

s=write_slide("🔴","火星","Mars",[
    ("火","huǒ","4 笔 — 跳动 的 火苗","Like flames"),
    ("星","xīng","9 笔 — 「日」+「生」","Sun + 生"),
],MARS); n+=1; pn(s,n)

# Session 3 divider
s=div("Session 3  下午 3:00–4:30","🛠️ 项目 · 太空 实验  ·  90 min",MARS,"🚀"); n+=1; pn(s,n)

# Projects overview
s=ns(); bg(s,CREAM); hb(s,"🛠️ 3 个 项目  3 Projects",MARS)
tb(s,0.4,0.85,9.2,0.32,"选 一个 你 最 想 做 的!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
projects=[
    ("🥤","吸管 火箭","Straw Rocket","纸 + 吸管 — 真的 飞!",MARS),
    ("🏠","火星 基地","Mars Base","设计 你 的 火星 家",NEBULA),
    ("💪","宇航员 训练","Astronaut Training","闭眼 / 平衡 / 慢动作",MOON_C),
]
for i,(em,cn,en,d,cl) in enumerate(projects):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.95),Inches(3.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,1.70,2.85,1.20,em,sz=78,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.85,0.42,cn,sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.38,2.85,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,2.75,0.85,d,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

# Project 1 — Straw rocket
s=ns(); bg(s,CREAM); hb(s,"🥤 项目 1 · 吸管 火箭  Straw Rocket",MARS)
ib(s,0.4,0.90,4.5,3.95,"🖼️ 吸管火箭 示例")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(3.95))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=MARS; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(0.90),Inches(4.50),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=MARS; head.line.fill.background()
tb(s,5.25,0.97,4.30,0.40,"📝 怎么 做",sz=14,b=True,c=WHITE)
steps=[
    ("1️⃣","用 一小张 纸 卷 成 一个 小管"),
    ("2️⃣","一头 用 胶带 封住 (火箭 头)"),
    ("3️⃣","剪 4 个 三角 — 当 火箭 尾翼"),
    ("4️⃣","贴 在 火箭 后面"),
    ("5️⃣","套 在 吸管 上 — 用力 吹! 🎯"),
]
for i,(num,txt) in enumerate(steps):
    y=1.55+i*0.62
    tb(s,5.25,y,0.40,0.40,num,sz=16,b=True,c=MARS)
    tb(s,5.75,y+0.04,3.80,0.35,txt,sz=11,c=DARK)
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.95),Inches(9.2),Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb=MARS; tip.line.fill.background()
tb(s,0.55,5.02,9.0,0.32,"💡 比一比 — 哪个 队 的 火箭 飞 得 最 远?",sz=12,b=True,c=STAR,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"项目 1 · 30 分钟:\n• 材料: 卡纸 + 吸管 + 胶带 + 剪刀\n• 课堂 在 外面 / 走廊 比赛 — 谁 飞 得 最 远")

# Projects 2 + 3 combined
s=ns(); bg(s,CREAM); hb(s,"🏠💪 项目 2 + 3",NEBULA)
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(4.10))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=NEBULA; left.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NEBULA; head.line.fill.background()
tb(s,0.55,0.97,4.30,0.40,"🏠 项目 2 · 火星 基地 设计",sz=14,b=True,c=WHITE)
items=[
    "🌟 火星 没有 空气 + 水 — 我们 要 自己 造!",
    "1️⃣ 画 你 的 火星 房子",
    "2️⃣ 必须 有: 氧气 室 + 水 + 食物",
    "3️⃣ 可以 加: 太阳能 板 / 温室 / 游乐场",
    "4️⃣ 用 彩笔 + 贴纸 装饰",
    "5️⃣ 给 你 的 基地 起 个 名字!",
]
for i,line in enumerate(items):
    tb(s,0.55,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(4.10))
right.fill.solid(); right.fill.fore_color.rgb=WHITE; right.line.color.rgb=MOON_C; right.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.90),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=MOON_C; head.line.fill.background()
tb(s,5.20,0.97,4.30,0.40,"💪 项目 3 · 宇航员 训练 挑战",sz=14,b=True,c=DARK)
parts=[
    "🌟 宇航员 训练 — 你 也 试 一试!",
    "1️⃣ 闭眼 走 5 步 (= 没有 GPS!)",
    "2️⃣ 单脚 站 30 秒 (= 月球 平衡)",
    "3️⃣ 慢动作 跳 10 下 (= 0 重力)",
    "4️⃣ 戴 手套 系 鞋带 (= 宇航服 太厚)",
    "5️⃣ 通过 → 拿 宇航员 徽章!",
]
for i,line in enumerate(parts):
    tb(s,5.20,1.55+i*0.42,4.30,0.40,line,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"30 分钟:\n• 项目 2 适合 喜欢 设计 的 学生\n• 项目 3 适合 喜欢 动 的 学生\n• 完成 → 每人 一个 宇航员 徽章 (sticker)")

# Share & close
s=ns(); bg(s,CREAM); hb(s,"🎤 分享 + 再见!",COSMIC)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(2.10))
sh.fill.solid(); sh.fill.fore_color.rgb=NIGHT; sh.line.color.rgb=STAR; sh.line.width=Pt(3)
tb(s,0.55,1.05,9.0,0.40,"💬 句型 — 分享 你 的 太空 任务!",sz=14,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,1.55,9.0,0.50,"「我 是 宇航员 — 我 去 ___」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
tb(s,0.55,2.10,9.0,0.50,"「我 看到 ___, 我 想 ___」",sz=22,b=True,c=STAR,a=PP_ALIGN.CENTER)
preview=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.25),Inches(9.2),Inches(2.05))
preview.fill.solid(); preview.fill.fore_color.rgb=WHITE; preview.line.color.rgb=COSMIC; preview.line.width=Pt(2.5)
tb(s,0.55,3.40,9.0,0.40,"🔮 下次 见 (Day 4):",sz=14,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
tb(s,0.55,3.85,9.0,0.40,"👽 外星人 + 来自 火星 的 信!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.30,9.0,0.30,"Aliens + a letter from Mars!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,4.70,9.0,0.30,"👋 你 觉得 外星人 真的 存在 吗? 想 一 想!",sz=11,b=True,c=COSMIC,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)

out=os.path.join(os.path.dirname(__file__),"day3_space_exploration.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
