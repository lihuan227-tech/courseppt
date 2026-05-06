#!/usr/bin/env python3
"""
野外生存与探险 Wilderness Unit — Day 5: 复习挑战日 / Review & Challenge Day
Capstone day — D1-2 review + Game 1, D3-4 review + Game 2, Final group challenge.
Modeled on little_artist_pbl/create_day5_exhibition.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Wilderness Adventure + Gold (capstone) ---
PINE     = RGBColor(0x1E,0x4D,0x2B)
SUN      = RGBColor(0xE0,0x7A,0x2C)
CREAM    = RGBColor(0xFD,0xF6,0xE3)
BROWN    = RGBColor(0x6B,0x44,0x23)
SKY      = RGBColor(0x4A,0x90,0xD9)
SUNYEL   = RGBColor(0xF5,0xC2,0x42)
ALERT    = RGBColor(0xD0,0x4A,0x3C)
GOLD     = RGBColor(0xFF,0xD7,0x00)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
DARK     = RGBColor(0x2C,0x2C,0x2C)
GRAY     = RGBColor(0x88,0x88,0x88)
LGRAY    = RGBColor(0xBB,0xBB,0xBB)
WARM     = RGBColor(0xFF,0xF3,0xE0)
IMGBG    = RGBColor(0xE8,0xE8,0xE8)
GREEN_OK = RGBColor(0x38,0x8E,0x3C)

# Per-day tags
FOREST   = RGBColor(0x2D,0x5A,0x3D)
RIVER    = RGBColor(0x19,0x76,0xD2)
NAVY     = RGBColor(0x1A,0x33,0x66)
DANGER   = RGBColor(0xC6,0x28,0x28)

PALETTE  = [PINE, SUN, BROWN, SKY, SUNYEL, ALERT, GOLD]

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=PINE,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def dot(s,x,y,r,color):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(r),Inches(r))
    d.fill.solid();d.fill.fore_color.rgb=color;d.line.fill.background()
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    other=[c for c in PALETTE if c!=color][:4]
    for (x,y),cl in zip([(0.8,4.7),(1.6,4.5),(7.8,4.5),(8.6,4.7)],other):
        dot(s,x,y,0.4,cl)
    return s

n=0

# ============================================================
# 1 COVER — Day 5 final badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.25,8,0.7,"Wilderness Adventure Camp",sz=32,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,1,0.85,8,0.45,"野外生存与探险夏令营",sz=20,c=PINE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.5),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=PINE;sh.line.color.rgb=GOLD;sh.line.width=Pt(8)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.8),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=GOLD;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.0,2.8,0.4,"DAY 5",sz=16,b=True,c=SUN,a=PP_ALIGN.CENTER)
ap(tf,"🏆",sz=58,a=PP_ALIGN.CENTER)
ap(tf,"复习挑战日",sz=20,b=True,c=PINE,a=PP_ALIGN.CENTER)
ap(tf,"REVIEW & CHALLENGE",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
for (x,y) in [(1.0,1.7),(8.6,1.6),(0.9,4.7),(8.7,4.6),(1.5,3.9),(8.0,3.9)]:
    dot(s,x,y,0.35,PALETTE[(int(x+y))%len(PALETTE)])
tb(s,1,5.05,8,0.4,"🎒 4 天的本领, 今天大比拼! Show what you've learned!",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE — three sessions
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule",SUN)
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","D1-D2 复习 + 团队比赛 Game 1",PINE),
    ("Session 2  下午","2:00-2:45","D3-D4 复习 + 团队比赛 Game 2",SUN),
    ("Session 3  下午","3:00-4:30","终极挑战 Final Challenge — 互相出题",GOLD)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    txt_color = DARK if cl==GOLD else WHITE
    sub_color = BROWN if cl==GOLD else WARM
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=txt_color)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=sub_color)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=txt_color)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives",SUN)
tb(s,0.5,0.9,9,0.5,"🌲 内容目标  Content:",sz=20,b=True,c=PINE)
tf=tb(s,0.7,1.4,9,1.6,"1. 复习 4 天所学：自然环境、营地、方向、生存安全",sz=15,c=DARK)
ap(tf,"2. 在团队比赛中应用所学知识",sz=15,c=DARK)
ap(tf,"3. 学会出题——巩固理解、与队友合作",sz=15,c=DARK)
ap(tf,"4. 在挑战中体验 探险家的勇气和自信",sz=15,c=DARK)
tb(s,0.5,3.2,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=SUN)
tb(s,0.7,3.7,9,0.4,"👀 综合复习: 森林 山地 河边 沙漠 雪地 营 帐篷 火 包 东南西北 指南针 天气 危险 迷路 食物 安全",sz=13,b=True,c=DARK)
tb(s,0.7,4.2,9,0.4,"✍️ 重点字: 森林 山地 河边 火 包 水 东 南 西 北 天气 危险",sz=13,b=True,c=DARK)
tb(s,0.5,4.75,9,0.5,"🎨 实践目标：完成 2 轮团队比赛 + 1 个终极挑战",sz=15,c=BROWN)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","D1 自然环境 + D2 营地搭建 复习\n🌲 6 种环境 · 🏕️ 营地选择 · 团队比赛 Game 1",PINE,"🔄")
n+=1

# ============================================================
# 5 D1 RECAP — 6 environments
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 Day 1 复习  6 种自然环境",PINE)
tb(s,0.4,0.9,9,0.35,"想一想 — 每种环境最大的危险是什么？",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
envs=[
    ("🌲","森林","Forest",FOREST,"迷路 / 动物"),
    ("🏔️","山地","Mountain",GRAY,"滑倒 / 天气"),
    ("🌾","草地","Grass",GREEN_OK,"晒 / 虫子"),
    ("🏞️","河边","River",RIVER,"溺水"),
    ("🏜️","沙漠","Desert",SUN,"中暑"),
    ("❄️","雪地","Snow",SKY,"冻伤"),
]
for i,(em,cn,en,cl,danger) in enumerate(envs):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.4+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.05,2.8,0.55,em,sz=32,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.7,2.8,0.4,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.05,2.8,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.32,2.8,0.3,f"⚠️ {danger}",sz=11,b=True,c=ALERT,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 6 D2 RECAP — Camping checklist + safe-site quick rules
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏕️ Day 2 复习  营地与装备",PINE)
# Left: gear cards
tb(s,0.4,0.85,4.5,0.4,"🎒 必带装备  Must-Pack Gear",sz=14,b=True,c=PINE)
gear=[
    ("⛺","帐篷 zhàng peng","tent"),
    ("🔥","火 huǒ","fire"),
    ("🎒","包 bāo","bag"),
    ("💧","水 shuǐ","water"),
    ("🍎","食物 shí wù","food"),
    ("🔦","手电筒","flashlight"),
]
for i,(em,cn,en) in enumerate(gear):
    col=i%2;row=i//2
    x=0.4+col*2.3;y=1.3+row*1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.2),Inches(1.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(1.5)
    tb(s,x+0.05,y+0.05,2.1,0.45,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.5,2.1,0.3,cn,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.78,2.1,0.25,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Right: safe-site rules
tb(s,5.2,0.85,4.5,0.4,"📍 营地选择  Safe Site Rules",sz=14,b=True,c=PINE)
rules=[
    ("✅ 平坦的地面","Flat ground"),
    ("✅ 远离河水 / 悬崖","Away from water / cliffs"),
    ("✅ 树下挡风, 不在大树正下方","Sheltered, not under one big tree"),
    ("❌ 不在低洼处 (会积水)","Not in dips (water collects)"),
    ("❌ 不在风口","Not in wind tunnels"),
]
for i,(cn,en) in enumerate(rules):
    y=1.3+i*0.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(y),Inches(4.5),Inches(0.55))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(1.2)
    tb(s,5.35,y+0.05,4.3,0.3,cn,sz=13,b=True,c=DARK)
    tb(s,5.35,y+0.3,4.3,0.25,en,sz=10,c=GRAY)
pn(s,n)

# ============================================================
# 7 GAME 1 INTRO — Team Competition
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 Game 1 · 团队大比拼  Team Competition",SUN)
tb(s,0.4,0.95,9.2,0.5,"分成 3-4 组, 每组取一个探险队名！",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.45,9.2,0.35,"Form 3-4 teams. Pick a wilderness team name!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# 3 round cards
rounds=[
    ("Round 1","🌲 环境抢答","Quick Quiz","看图说环境 + 危险",PINE),
    ("Round 2","🏕️ 营地判断","Site Judge","哪里能搭? 哪里不能?",SUN),
    ("Round 3","🎯 我会认","Word Match","D1+D2 字卡配对",BROWN),
]
for i,(num,cn,en,desc,cl) in enumerate(rounds):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(3.0),Inches(2.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.05,2.8,0.4,num,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.5,2.8,0.5,cn,sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.05,2.8,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.5,2.8,0.9,desc,sz=12,c=DARK,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.75),Inches(9.4),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2)
tb(s,0.5,4.83,9.0,0.4,"🏆 答对 +1 分, 用中文加 +1 额外分! Correct +1, in Chinese +1 bonus!",sz=13,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8 GAME 1 ROUND 1 — Environment quick quiz
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌲 Round 1 · 环境抢答  Environment Quiz",PINE)
tb(s,0.4,0.85,9.2,0.35,"老师出图, 哪组先抢到说出 中文 + 危险，就得分",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("🌲","森林","迷路 / 看不清路"),
    ("🌾","草地","虫子 / 太阳晒"),
    ("🏞️","河边","溺水 / 滑倒"),
    ("❄️","雪地","冻伤 / 滑"),
    ("🏔️","山地","滑倒 / 天气变"),
    ("🏜️","沙漠","中暑 / 缺水"),
]
for i,(em,cn,danger) in enumerate(qs):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.3+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.1,2.8,0.7,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.95,2.8,0.4,cn,sz=18,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.4,2.8,0.3,f"⚠️ {danger}",sz=11,c=ALERT,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 9 GAME 1 ROUND 2 — Site judge
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏕️ Round 2 · 营地判断  Site Judge",SUN)
tb(s,0.4,0.85,9.2,0.35,"看图说: 这里能搭帐篷吗? 为什么? Can we camp here? Why?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
sites=[
    ("🌳","树下平地","✅ 好",GREEN_OK,"挡风, 平坦"),
    ("🏞️","河边湿地","❌ 不",ALERT,"会涨水, 滑"),
    ("🏔️","山顶","❌ 不",ALERT,"风大, 雷击"),
    ("🌾","草地","✅ 好",GREEN_OK,"平坦, 远离水"),
    ("🕳️","低洼坑","❌ 不",ALERT,"下雨会积水"),
    ("🌬️","风口","❌ 不",ALERT,"风太大"),
]
for i,(em,site,ok,cl,reason) in enumerate(sites):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.3+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.05,2.8,0.5,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.62,2.8,0.4,site,sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.0,2.8,0.4,ok,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.4,2.8,0.3,reason,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 10 GAME 1 ROUND 3 — Word match (D1+D2 chars)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 Round 3 · 我会认 字卡配对  Word Match",BROWN)
tb(s,0.4,0.85,9.2,0.35,"老师举字卡, 各组抢答中文+英文; 用中文造句 +1 额外分",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
words=[
    ("森林","sēn lín","Forest","🌲"),
    ("山地","shān dì","Mountain","🏔️"),
    ("河边","hé biān","Riverside","🏞️"),
    ("帐篷","zhàng peng","Tent","⛺"),
    ("火","huǒ","Fire","🔥"),
    ("包","bāo","Bag","🎒"),
    ("营","yíng","Camp","🏕️"),
    ("安全","ān quán","Safe","✅"),
]
for i,(cn,py,en,em) in enumerate(words):
    col=i%4;row=i//4
    x=0.3+col*2.4;y=1.3+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.25),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BROWN;sh.line.width=Pt(2)
    tb(s,x+0.05,y+0.05,2.15,0.45,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.55,2.15,0.45,cn,sz=22,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.05,2.15,0.3,py,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.32,2.15,0.3,en,sz=11,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 11 SCOREBOARD 1
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 Game 1 计分板  Scoreboard",GOLD)
tb(s,0.4,0.95,9.2,0.4,"老师在白板上画计分表: 哪个队中文说得多, 哪个队就赢!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9.2,0.35,"Teacher tallies on whiteboard — Chinese answers count more!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Mock scoreboard
teams=["Team 1","Team 2","Team 3","Team 4"]
for i,tm in enumerate(teams):
    x=0.5+i*2.3
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(2.1),Inches(2.8))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PALETTE[i%len(PALETTE)];sh.line.width=Pt(3)
    tb(s,x+0.1,2.05,1.9,0.4,tm,sz=14,b=True,c=PALETTE[i%len(PALETTE)],a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.5,1.9,0.4,"队名:",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.85,1.9,0.4,"_________",sz=14,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.5,1.9,0.4,"得分 Score:",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.9,1.9,0.7,"____",sz=36,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,5.0,9.2,0.35,"🌟 答对 +1 · 中文回答 +1 额外分 · 队伍合作 +1",sz=12,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 12 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","D3 方向地图 + D4 生存安全 复习\n🧭 东南西北 · ⚠️ 紧急情况 · 团队比赛 Game 2",SUN,"🔄")
n+=1

# ============================================================
# 13 D3 RECAP — Compass + 4 directions
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧭 Day 3 复习  方向与指南针",SUN)
# Big compass center
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=NAVY;sh.line.width=Pt(4)
tb(s,3.5,2.2,3,1.4,"🧭",sz=80,a=PP_ALIGN.CENTER)
tb(s,3.5,3.6,3,0.4,"指南针",sz=18,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,3.5,3.95,3,0.3,"zhǐ nán zhēn · Compass",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 4 direction markers
dirs=[("北 North","běi","❄️ 冷",4.7,0.95),("南 South","nán","🥵 热",4.7,4.5),
      ("东 East","dōng","🌅 日出",1.0,2.6),("西 West","xī","🌇 日落",7.5,2.6)]
for cn,py,hint,x,y in dirs:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x-0.3),Inches(y),Inches(2.0),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=NAVY;sh.line.width=Pt(2)
    tb(s,x-0.25,y+0.05,1.9,0.4,cn,sz=15,b=True,c=NAVY,a=PP_ALIGN.CENTER)
    tb(s,x-0.25,y+0.45,1.9,0.3,f"{py} · {hint}",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.95,9.2,0.35,"💡 提示: 太阳早上从东边升起, 晚上从西边落下",sz=12,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 14 D4 RECAP — Emergency situations
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⚠️ Day 4 复习  生存与安全",DANGER)
tb(s,0.4,0.85,9.2,0.35,"遇到紧急情况, 你应该怎么做？  Emergency response",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
emerg=[
    ("🌧️","坏天气","找掩护, 不在树下"),
    ("🗺️","迷路","原地不动, 找大人"),
    ("🐻","野生动物","安静慢慢离开"),
    ("🍎","食物","不认识的不能吃"),
    ("💧","水","只喝瓶装水或煮开"),
    ("📞","求救","SOS 三声三停"),
]
for i,(em,sit,act) in enumerate(emerg):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.3+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.05,2.8,0.5,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.6,2.8,0.4,sit,sz=15,b=True,c=DANGER,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.05,2.8,0.5,act,sz=11,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 15 GAME 2 INTRO
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 Game 2 · 团队大比拼 第二轮",SUN)
tb(s,0.4,0.95,9.2,0.5,"换组! 重新分组也可以, 探险家要适应新队友",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.45,9.2,0.35,"Reshuffle if you like — explorers adapt to new teammates!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
rounds2=[
    ("Round 1","🧭 方向闪问","Direction Flash","看图说东南西北",NAVY),
    ("Round 2","⚠️ 危险情景","Danger Scene","看场景说怎么办",DANGER),
    ("Round 3","🎯 我会认","Word Match","D3+D4 字卡配对",BROWN),
]
for i,(num,cn,en,desc,cl) in enumerate(rounds2):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(3.0),Inches(2.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.05,2.8,0.4,num,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.5,2.8,0.5,cn,sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.05,2.8,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.5,2.8,0.9,desc,sz=12,c=DARK,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.75),Inches(9.4),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2)
tb(s,0.5,4.83,9.0,0.4,"🏆 同样规则: 答对 +1, 中文回答 +1 额外分",sz=13,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 16 GAME 2 ROUND 1 — Direction quiz
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧭 Round 1 · 方向闪问  Direction Flash",NAVY)
tb(s,0.4,0.85,9.2,0.35,"老师给情景, 各组抢答方向 — 用中文 +1",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("☀️","太阳从哪边升起？","东 East"),
    ("🌇","太阳从哪边落下？","西 West"),
    ("❄️","哪个方向最冷？","北 North"),
    ("🥵","哪个方向最热？","南 South"),
    ("🧭","红针指哪边？","北 North"),
    ("🌳","北半球苔藓长在哪边？","北 North"),
]
for i,(em,q,ans) in enumerate(qs):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.3+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=NAVY;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.05,2.8,0.5,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.6,2.8,0.5,q,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.15,2.8,0.4,f"答: {ans}",sz=13,b=True,c=NAVY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 17 GAME 2 ROUND 2 — Danger scene
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⚠️ Round 2 · 危险情景  Danger Scene",DANGER)
tb(s,0.4,0.85,9.2,0.35,"老师描述场景, 各组抢答怎么做 — 用中文 +1",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
scenes=[
    ("🌧️⛈️","下大雨打雷","找掩护, 不要在树下"),
    ("🐻","看到一只熊","安静慢慢离开"),
    ("🗺️❓","和大家走散了","原地等, 大声呼救"),
    ("🍄","看到漂亮蘑菇","不认识的不能吃"),
    ("🌊","脚下湿石头","小心慢走或绕开"),
    ("🥵","太阳很大","戴帽子, 多喝水"),
]
for i,(em,sit,act) in enumerate(scenes):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.3+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.05,2.8,0.5,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.6,2.8,0.4,sit,sz=14,b=True,c=DANGER,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.05,2.8,0.5,f"✅ {act}",sz=11,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 18 GAME 2 ROUND 3 — D3+D4 word match
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 Round 3 · 我会认 字卡配对",BROWN)
tb(s,0.4,0.85,9.2,0.35,"D3 + D4 字卡 — 老师举牌, 各组抢答",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
words2=[
    ("东","dōng","East","🌅"),
    ("南","nán","South","🥵"),
    ("西","xī","West","🌇"),
    ("北","běi","North","❄️"),
    ("指南针","zhǐ nán zhēn","Compass","🧭"),
    ("天气","tiān qì","Weather","🌤️"),
    ("危险","wēi xiǎn","Danger","⚠️"),
    ("迷路","mí lù","Lost","🗺️"),
]
for i,(cn,py,en,em) in enumerate(words2):
    col=i%4;row=i//4
    x=0.3+col*2.4;y=1.3+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.25),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BROWN;sh.line.width=Pt(2)
    tb(s,x+0.05,y+0.05,2.15,0.45,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.55,2.15,0.45,cn,sz=20,b=True,c=NAVY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.05,2.15,0.3,py,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.32,2.15,0.3,en,sz=11,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 19 SCOREBOARD 2 — running total
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 Game 2 计分板 + 总分  Total Scoreboard",GOLD)
tb(s,0.4,0.95,9.2,0.4,"把 Game 1 + Game 2 加起来 — 谁是探险大王?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9.2,0.35,"Game 1 + Game 2 totals — who's the Wilderness Champion?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
rows=[("Game 1 Score","____","____","____","____"),
      ("Game 2 Score","____","____","____","____"),
      ("Total 总分","____","____","____","____")]
# Header
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.4),Inches(2.0),Inches(2.0),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.fill.background()
tb(s,0.4,2.05,2.0,0.4,"",sz=12)
for i,tm in enumerate(["Team 1","Team 2","Team 3","Team 4"]):
    x=2.4+i*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.0),Inches(1.8),Inches(0.5))
    sh.fill.solid();sh.fill.fore_color.rgb=PALETTE[i%len(PALETTE)];sh.line.fill.background()
    tb(s,x,2.05,1.8,0.4,tm,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
for r,row in enumerate(rows):
    y=2.5+r*0.7
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.4),Inches(y),Inches(2.0),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=GRAY;sh.line.width=Pt(0.5)
    tb(s,0.5,y+0.15,1.9,0.4,row[0],sz=12,b=True,c=DARK)
    for i in range(4):
        x=2.4+i*1.85
        sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(1.8),Inches(0.65))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GRAY;sh.line.width=Pt(0.5)
        tb(s,x,y+0.1,1.8,0.5,row[i+1],sz=18,b=(r==2),c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 20 SESSION 3 DIVIDER — Final Challenge (GOLD)
# ============================================================
div("Session 3  下午","🏆 终极挑战 Final Challenge\n各组互相出 10 道题 — 看谁能闯关成功!",GOLD,"🎯")
n+=1

# ============================================================
# 21 FINAL CHALLENGE OVERVIEW
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 终极挑战  Final Challenge — 互相出题",GOLD)
tb(s,0.4,0.85,9.2,0.4,"每组给另一组 10 道题, 看对方能不能闯关!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.3,9.2,0.35,"Each team writes 10 challenges for another team — can they pass?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("1️⃣","分组出题","20 分钟","每组讨论, 出 10 道题"),
    ("2️⃣","抽签换题","2 分钟","抽到哪组就答哪组的题"),
    ("3️⃣","闯关答题","20 分钟","按顺序回答 10 道题"),
    ("4️⃣","计分公布","8 分钟","看谁分最高 = 探险大王"),
]
for i,(num,t,tm,desc) in enumerate(steps):
    col=i%2;row=i//2
    x=0.4+col*4.7;y=1.85+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(1.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GOLD;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.1,0.7,0.5,num,sz=24,b=True,c=GOLD,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.1,3.5,0.4,t,sz=15,b=True,c=DARK)
    tb(s,x+0.85,y+0.5,3.5,0.3,f"⏱️ {tm}",sz=11,c=SUN)
    tb(s,x+0.85,y+0.85,3.5,0.4,desc,sz=11,c=DARK)
pn(s,n)

# ============================================================
# 22 STEP 1 — Write 10 questions
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"1️⃣ 分组出题 (20 分钟)  Write 10 Questions",PINE)
tb(s,0.4,0.9,9.2,0.4,"每组要出 10 道题, 题目要覆盖 4 天的内容",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Left: question types
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.45),Inches(4.6),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(2.5)
tb(s,0.5,1.55,4.4,0.35,"📝 题目类型 Question Types",sz=14,b=True,c=PINE)
qtypes=[
    ("看图选答案","Pick from image"),
    ("说出中文 / 英文","Say in CN / EN"),
    ("说出特点 / 危险","Tell features / dangers"),
    ("情景判断","Scene judgment"),
    ("写一个字","Write a character"),
]
for i,(cn,en) in enumerate(qtypes):
    y=2.0+i*0.55
    tb(s,0.5,y,4.4,0.3,f"• {cn}",sz=13,b=True,c=DARK)
    tb(s,0.5,y+0.28,4.4,0.25,en,sz=10,c=GRAY)
# Right: tip
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.45),Inches(4.6),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
tb(s,5.25,1.55,4.4,0.35,"💡 出题小建议  Tips",sz=14,b=True,c=SUN)
tips=[
    ("覆盖 4 天","Cover all 4 days"),
    ("不能用问问题让队伍输, 要公平","Be fair, not too tricky"),
    ("能用中文写就用中文写","Write in Chinese when you can"),
    ("准备好 答案 + 加分理由","Prep answers + bonus rules"),
]
for i,(cn,en) in enumerate(tips):
    y=2.0+i*0.65
    tb(s,5.25,y,4.4,0.3,f"✓ {cn}",sz=13,b=True,c=DARK)
    tb(s,5.25,y+0.28,4.4,0.25,en,sz=10,c=GRAY)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.05),Inches(9.4),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=PINE;sh.line.fill.background()
tb(s,0.5,5.13,9.0,0.4,"⏱️ 20 分钟内完成 — 团队合作!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 23 STEP 2 — Draw lots
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"2️⃣ 抽签换题 (2 分钟)  Draw Lots",SUN)
tb(s,0.4,0.95,9.2,0.4,"每组派一人抽签 — 抽到谁的题, 就答谁的题",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9.2,0.35,"One member from each team draws — answer the questions from that team",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Visual: 4 hat / lottery icons
for i,em in enumerate(["🎩","✋","📜","🎯"]):
    x=0.5+i*2.3
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(2.1),Inches(2.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
    tb(s,x+0.1,2.1,1.9,1.0,em,sz=64,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.3,1.9,0.4,f"Team {i+1}",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.7,1.9,0.4,"答 → ?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.05,1.9,0.4,"_______",sz=14,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.85,9.2,0.35,"📌 不能抽到自己组的题 — 让老师重抽",sz=12,b=True,c=ALERT,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 24 STEP 3 — Run challenges (rules)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"3️⃣ 闯关答题 (20 分钟)  Run the Challenges",ALERT)
tb(s,0.4,0.9,9.2,0.4,"⚖️ 规则  Rules",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
# 4 rule cards
rules=[
    ("✅","答对",f"答题组 +1\n(出题组 0)",GREEN_OK),
    ("❌","答错",f"出题组 +1\n(答题组 0)",ALERT),
    ("🇨🇳","用中文答",f"答题组 +1\n额外加分",SUN),
    ("🤝","团队合作",f"全组讨论后\n再回答",PINE),
]
for i,(em,t,desc,cl) in enumerate(rules):
    col=i%2;row=i//2
    x=0.4+col*4.7;y=1.5+row*1.7
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(1.55))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.1,0.7,0.6,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.1,3.5,0.4,t,sz=16,b=True,c=cl)
    tb(s,x+0.85,y+0.55,3.5,0.95,desc,sz=12,c=DARK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.95),Inches(9.4),Inches(0.55))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.fill.background()
tb(s,0.5,5.05,9.0,0.4,"💎 中文写得越多, 队伍分数越高! More Chinese = more points!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 25 STEP 4 — Final Scoreboard
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"4️⃣ 公布最终分  Final Scoreboard",GOLD)
tb(s,0.4,0.9,9.2,0.4,"把今天所有的分加起来 — 探险大王诞生!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Big trophy
tb(s,4,1.4,2,1.4,"🏆",sz=110,a=PP_ALIGN.CENTER)
# 3 podium boxes
podium=[("🥇","第一名","Champion",2.5),("🥈","第二名","Runner-up",1.7),("🥉","第三名","3rd Place",0.9)]
for i,(em,cn,en,h) in enumerate(podium):
    x=2.0+i*2.0
    y=4.0-h*0.4
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(1.8),Inches(h*0.4+0.6))
    sh.fill.solid();sh.fill.fore_color.rgb=[GOLD,LGRAY,RGBColor(0xCD,0x7F,0x32)][i];sh.line.fill.background()
    tb(s,x,y+0.05,1.8,0.4,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.5,1.8,0.4,cn,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.85,1.8,0.3,en,sz=10,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,5.0,9.2,0.4,"🎉 每组都有探险家徽章 — 大家都是赢家! Everyone earns a badge!",sz=13,b=True,c=PINE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 26 BADGE — Day 5 + Unit complete
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.4,9,0.8,"🎖️ Day 5 探险家徽章  Wilderness Master Badge",sz=26,b=True,c=PINE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GOLD;sh.line.width=Pt(6)
tf=tb(s,3.6,1.65,2.8,2.7,"DAY 5",sz=18,b=True,c=SUN,a=PP_ALIGN.CENTER)
ap(tf,"🏆",sz=42,a=PP_ALIGN.CENTER)
ap(tf,"复习挑战日",sz=18,b=True,c=PINE,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=12,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
ap(tf,"🌲🏕️🧭⚠️🏆",sz=14,a=PP_ALIGN.CENTER)
tb(s,1,4.55,8,0.4,"🎉 恭喜你完成 Wilderness Unit！ Congratulations, young explorer!",sz=15,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"4 天本领 · 6 种环境 · 4 个方向 · 无数勇气",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 27 CLOSING — Thanks + farewell
# ============================================================
s=ns();n+=1;bg(s,PINE)
tb(s,0.5,1.0,9,0.8,"🌲 感谢你的旅程!  Thank you, explorers!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.4,7,2.5,"Wilderness Adventure 完美收官",sz=24,b=True,c=GOLD,a=PP_ALIGN.CENTER)
ap(tf,"Camp completed!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🎒 带着勇气和知识, 走向更大的世界",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Take your courage and skills into the wider world",sz=12,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"下次再见, 小探险家! See you next time!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/huanli/projects/courseppt/Chinese/野外生存与探险wilderness_pbl/day5_review.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
