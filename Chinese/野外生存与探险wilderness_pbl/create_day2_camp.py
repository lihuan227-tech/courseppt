#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
野外生存与探险 Day 2 — 搭建营地 Build a Mini Camp (Camping & Safety)
"Explorer Mission" framing — kids = small explorers building a safe camp.
Each section = a task: pack, build, choose location, protect.
Reuses Day 1 palette + helper conventions.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette (continuity with Day 1) ---
PINE   = RGBColor(0x1E,0x4D,0x2B)
SUN    = RGBColor(0xE0,0x7A,0x2C)
CREAM  = RGBColor(0xFD,0xF6,0xE3)
BROWN  = RGBColor(0x6B,0x44,0x23)
SKY    = RGBColor(0x4A,0x90,0xD9)
SUNYEL = RGBColor(0xF5,0xC2,0x42)
ALERT  = RGBColor(0xD0,0x4A,0x3C)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
WARM   = RGBColor(0xFF,0xF3,0xE0)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
OK     = RGBColor(0x38,0x8E,0x3C)
# Camp-zone accents
TENTCL = RGBColor(0x2D,0x5A,0x3D)
FIRECL = RGBColor(0xE6,0x52,0x2C)
FOODCL = RGBColor(0xC9,0x82,0x46)
SAFCL  = RGBColor(0x3E,0x6E,0xB6)
PLAYCL = RGBColor(0x8E,0x44,0xAD)  # purple — recreation/play zone

# === Helpers (same conventions as Day 1) ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=PINE,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,txt):
    s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER);return s
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)

# === Specialized Day 2 helpers ===

def mission_card(s,l,t,w,h,num,task_cn,task_en,emoji,color):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l+0.1),Inches(t+0.08),Inches(0.55),Inches(0.55))
    badge.fill.solid();badge.fill.fore_color.rgb=color;badge.line.fill.background()
    tb(s,l+0.1,t+0.18,0.55,0.4,str(num),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,l+0.7,t+0.15,w-0.8,0.4,task_en,sz=10,c=GRAY)
    tb(s,l+0.05,t+0.85,w-0.1,0.7,emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+1.55,w-0.1,0.4,task_cn,sz=18,b=True,c=color,a=PP_ALIGN.CENTER)

def sentence_frame_bar(s,t,frame_cn,frame_en):
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.7,0.4,"💬 我来说",sz=14,b=True,c=SUN)
    tb(s,2.0,t+0.07,7.6,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.0,t+0.32,7.6,0.3,frame_en,sz=10,c=GRAY,a=None)

def zone_slide(emoji,name_cn,name_en,color,rules,frame_cn,frame_en):
    """A camp-zone slide: 4 simple safety rules + sentence frame + image placeholder."""
    s=ns();bg(s,CREAM);hb(s,f"{emoji} 营地 4 区域 · {name_cn} {name_en}",color)
    # left: image
    ib(s,0.3,1.0,4.3,3.3,f"📷 {name_cn} 图片 / 简笔画")
    # right: rule cards
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.0),Inches(4.85),Inches(3.3))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=color;panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.0),Inches(4.85),Inches(0.5))
    head.fill.solid();head.fill.fore_color.rgb=color;head.line.fill.background()
    tb(s,5.0,1.07,4.6,0.4,"🤔 安全的 / 不安全的？",sz=14,b=True,c=WHITE)
    y=1.62
    for icon,rule in rules:
        tb(s,5.05,y,0.4,0.35,icon,sz=18)
        tb(s,5.45,y+0.02,4.0,0.35,rule,sz=13,c=DARK);y+=0.6
    sentence_frame_bar(s,4.45,frame_cn,frame_en)
    return s

def zone_q_slide(emoji,zone_cn,zone_en,color,questions,frame_cn,frame_en,img_label="📷 营区图片"):
    """Per-zone QUESTION slide — image LEFT + questions RIGHT, students think first."""
    s=ns();bg(s,CREAM);hb(s,f"{emoji} {zone_cn} · 🤔 想一想 Think First",color)
    tb(s,0.4,0.80,9.2,0.30,"先看图想一想 — 你怎么想? Look & think before tips!",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # LEFT — image area (paste real picture later)
    img_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.20),Inches(4.30),Inches(3.20))
    img_box.fill.solid();img_box.fill.fore_color.rgb=IMGBG;img_box.line.color.rgb=color;img_box.line.width=Pt(2)
    tb(s,0.5,2.50,4.10,0.40,img_label,sz=12,b=True,c=LGRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.90,4.10,0.30,zone_en,sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
    # RIGHT — questions stacked
    y0=1.20
    for i,(q_cn,q_en) in enumerate(questions[:3]):
        y=y0+i*1.05
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(y),Inches(4.85),Inches(0.95))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(4.98),Inches(y+0.22),Inches(0.50),Inches(0.50))
        badge.fill.solid();badge.fill.fore_color.rgb=color;badge.line.fill.background()
        tb(s,4.98,y+0.27,0.50,0.4,str(i+1),sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,5.60,y+0.12,4.05,0.40,q_cn,sz=14,b=True,c=DARK)
        tb(s,5.60,y+0.52,4.05,0.32,q_en,sz=9,c=GRAY)
    sentence_frame_bar(s,4.55,frame_cn,frame_en)
    return s

def answer_panels_slide(header_text, color, panels, img_label="📷 营区图片", img_path=None, subtitle="揭晓答案 — 看图记住!  Now the tips — see the picture, remember!"):
    """Image-format answer slide: photo LEFT, ✅/❌ panels with mixed CN+EN text RIGHT.
    - header_text: full header (e.g. "🏕️ 帐篷区 · 💡 露营小贴士 Pro Tips")
    - panels: list of dicts {"q": "Question?", "mark": "✅" or "❌", "lines": [mixed CN+EN strings]}
      "q" is optional — if present, shown as panel-color question header above the bullets.
    - img_path: optional real image path; if missing, uses placeholder
    """
    s=ns();bg(s,CREAM);hb(s,header_text,color)
    tb(s,0.4,0.74,9.2,0.26,subtitle,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # LEFT — image area
    img_top=1.04; img_h=4.30
    img_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.30),Inches(img_top),Inches(4.30),Inches(img_h))
    img_box.fill.solid();img_box.fill.fore_color.rgb=IMGBG;img_box.line.color.rgb=color;img_box.line.width=Pt(2)
    if img_path and os.path.exists(img_path):
        s.shapes.add_picture(img_path,Inches(0.40),Inches(img_top+0.10),Inches(4.10),Inches(img_h-0.20))
    else:
        tb(s,0.30,img_top+img_h/2-0.20,4.30,0.40,img_label,sz=12,b=True,c=LGRAY,a=PP_ALIGN.CENTER)
    # RIGHT — answer panels (same vertical bounds as image)
    px=4.80; pw=4.90
    n_panels=len(panels)
    total_avail=img_h
    gap=0.10
    each_h=(total_avail - gap*(n_panels-1))/n_panels
    y=img_top
    for p in panels:
        mark=p["mark"]; lines=p["lines"]; q=p.get("q")
        mark_color = OK if mark=="✅" else ALERT
        # Panel box
        pb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(px),Inches(y),Inches(pw),Inches(each_h))
        pb.fill.solid();pb.fill.fore_color.rgb=WHITE
        pb.line.color.rgb=mark_color;pb.line.width=Pt(2.5)
        # Mark on left (centered vertically)
        tb(s,px+0.08,y+each_h/2-0.25,0.55,0.50,mark,sz=22,b=True,c=mark_color,a=PP_ALIGN.CENTER)
        # Text area: question header (panel-color) + answer bullets (DARK)
        text_x=px+0.72; text_w=pw-0.82
        line_y=y+0.07
        bx=s.shapes.add_textbox(Inches(text_x),Inches(line_y),Inches(text_w),Inches(each_h-0.14))
        tf=bx.text_frame; tf.word_wrap=True
        tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
        tf.margin_left=Pt(0); tf.margin_right=Pt(0)
        para_idx=0
        if q:
            # Question header
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            try:
                para.space_before = Pt(0)
                para.space_after = Pt(1)
                para.line_spacing = 1.0
            except Exception:
                pass
            r = para.add_run()
            r.text = f"❓ {q}"
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = mark_color
            r.font.name = 'KaiTi'
            para_idx = 1
        for i,ln in enumerate(lines):
            if i==0 and para_idx==0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            try:
                para.space_before = Pt(1)
                para.space_after = Pt(1)
                para.line_spacing = 1.0
            except Exception:
                pass
            r = para.add_run()
            r.text = ln
            r.font.size = Pt(10)
            r.font.bold = True
            r.font.color.rgb = DARK
            r.font.name = 'KaiTi'
        y += each_h + gap
    return s

def zone_tips_slide(emoji,zone_cn,zone_en,color,tips,frame_cn,frame_en,img_label="📷 营区图片"):
    """Per-zone TIPS reveal slide — image LEFT + tip list RIGHT."""
    s=ns();bg(s,CREAM);hb(s,f"{emoji} {zone_cn} · 💡 露营小贴士 Pro Tips",color)
    tb(s,0.4,0.80,9.2,0.30,"揭晓答案 — 看图记住!  Now the tips — see the picture, remember!",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # LEFT — image area
    img_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.20),Inches(4.30),Inches(3.20))
    img_box.fill.solid();img_box.fill.fore_color.rgb=IMGBG;img_box.line.color.rgb=color;img_box.line.width=Pt(2)
    tb(s,0.5,2.50,4.10,0.40,img_label,sz=12,b=True,c=LGRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.90,4.10,0.30,zone_en,sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
    # RIGHT — compact tip list (max 4 tips fit nicely)
    tips_show=tips[:4]
    y0=1.20
    h_each=(3.20-(len(tips_show)-1)*0.10)/len(tips_show)
    for i,(mark,tip_cn,tip_en) in enumerate(tips_show):
        y=y0+i*(h_each+0.10)
        mark_color = OK if mark=="✅" else ALERT
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(y),Inches(4.85),Inches(h_each))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(1.5)
        tb(s,4.95,y+(h_each-0.4)/2,0.45,0.4,mark,sz=18,b=True,c=mark_color,a=PP_ALIGN.CENTER)
        tb(s,5.45,y+0.08,4.20,0.35,tip_cn,sz=13,b=True,c=DARK)
        tb(s,5.45,y+0.42,4.20,0.30,tip_en,sz=9,c=GRAY)
    sentence_frame_bar(s,4.55,frame_cn,frame_en)
    return s

def ab_slide(title_cn,title_en,question_cn,question_en,a_emoji,a_label,a_caption,b_emoji,b_label,b_caption,answer,reason):
    """A vs B vote slide — students physically move left/right."""
    s=ns();bg(s,CREAM);hb(s,f"📍 {title_cn}  {title_en}",SAFCL)
    # question
    tb(s,0.4,0.85,9.2,0.4,question_cn,sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.20,9.2,0.3,question_en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    # A panel (left)
    a_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.6),Inches(4.5),Inches(2.4))
    a_box.fill.solid();a_box.fill.fore_color.rgb=WHITE;a_box.line.color.rgb=SAFCL;a_box.line.width=Pt(2.5)
    pill(s,0.5,1.7,0.7,0.4,"A",SAFCL,sz=16)
    tb(s,1.2,1.65,3.6,0.5,a_label,sz=18,b=True,c=SAFCL)
    tb(s,0.6,2.2,4.2,0.5,a_emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,0.6,3.1,4.2,0.6,a_caption,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    # B panel (right)
    b_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.6),Inches(4.5),Inches(2.4))
    b_box.fill.solid();b_box.fill.fore_color.rgb=WHITE;b_box.line.color.rgb=ALERT;b_box.line.width=Pt(2.5)
    pill(s,5.2,1.7,0.7,0.4,"B",ALERT,sz=16)
    tb(s,5.9,1.65,3.6,0.5,b_label,sz=18,b=True,c=ALERT)
    tb(s,5.3,2.2,4.2,0.5,b_emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,5.3,3.1,4.2,0.6,b_caption,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    # vote prompt
    sentence_frame_bar(s,4.15,
        "我选 ___, 因为 ____  /  我应该 ____",
        "I choose A/B because… / I should…")
    # Footer hint
    tb(s,0.4,4.92,9.2,0.3,"👉 走到 A 边 / B 边 — 用一句话说为什么。",sz=12,b=True,c=SUN,a=PP_ALIGN.CENTER)
    notes(s,f"老师备课:\n• 答案 / Answer: {answer}\n• 原因 / Reason: {reason}\n• 玩法: 把教室分成 A/B 两边, 学生走到选择的一边, 然后每边请 1-2 个孩子说出原因。\n• K 级: 「A 安全」「B 不安全」即可。\n• G1-3: 「我选 ___, 因为 ___」整句。\n• 教师可在视频/图片资源里截图替换 emoji。")
    return s

def video_slide(title_cn,title_en,before_task,after_action,bgc=SAFCL):
    """Video slide with before-watching task + after-watching action."""
    s=ns();bg(s,bgc)
    tb(s,1,0.55,8,0.7,"🎬 看视频 Watch Video",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,1.25,8,0.4,f"{title_cn}",sz=20,b=True,c=WARM,a=PP_ALIGN.CENTER)
    tb(s,1,1.6,8,0.3,title_en,sz=12,c=WARM,a=PP_ALIGN.CENTER)
    # Before
    pre=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(2.05),Inches(4.5),Inches(2.0))
    pre.fill.solid();pre.fill.fore_color.rgb=WHITE;pre.line.fill.background()
    tb(s,0.55,2.15,4.3,0.4,"👂 看之前 Before Watching",sz=14,b=True,c=bgc)
    tb(s,0.55,2.55,4.3,1.4,before_task,sz=12,c=DARK)
    # After
    post=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(2.05),Inches(4.5),Inches(2.0))
    post.fill.solid();post.fill.fore_color.rgb=WHITE;post.line.fill.background()
    tb(s,5.25,2.15,4.3,0.4,"🎯 看完后 After Watching",sz=14,b=True,c=SUN)
    tb(s,5.25,2.55,4.3,1.4,after_action,sz=12,c=DARK)
    # Link placeholder
    tb(s,1,4.3,8,0.4,"🔗 视频链接 (老师粘贴) / Teacher pastes video link",sz=12,c=WARM,a=PP_ALIGN.CENTER)
    return s

def word_card_read(w,py,en,sent,img,color=SUN):
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=color,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句 Example",sz=14,b=True,c=color)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=20,b=True,c=DARK)
    return s

def word_card_write(w,py,en,strokes_hint,color=PINE):
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    # Stroke order area
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=color)
    tb(s,0.6,3.8,4.6,1.2,strokes_hint,sz=14,c=DARK)
    # Practice steps
    tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=color)
    tb(s,5.8,3.8,3.8,1.2,"1. 看老师写\n2. 用手指空中写\n3. 在本子上写 3 次",sz=13,c=DARK)
    # Right: 田字格 boxes (4 squares)
    for i in range(4):
        sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(5.8+i*0.95),Inches(1.0),Inches(0.85),Inches(0.85))
        sq.fill.solid();sq.fill.fore_color.rgb=WHITE;sq.line.color.rgb=color;sq.line.width=Pt(1.5)
        # cross hair
        ln1=s.shapes.add_connector(1,Inches(5.8+i*0.95),Inches(1.425),Inches(5.8+i*0.95+0.85),Inches(1.425))
        ln1.line.color.rgb=LGRAY;ln1.line.width=Pt(0.5);ln1.line.dash_style=2
        ln2=s.shapes.add_connector(1,Inches(5.8+i*0.95+0.425),Inches(1.0),Inches(5.8+i*0.95+0.425),Inches(1.85))
        ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.5);ln2.line.dash_style=2
    tb(s,5.8,1.95,3.8,0.3,"在田字格里写 3 次 ↓",sz=11,c=GRAY)
    return s

# ========================================================================
#                              SLIDES
# ========================================================================
n=0

# 1. COVER
s=ns();bg(s,PINE)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.4,8,0.5,"DAY 2",sz=18,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"🏕️ 搭建营地",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"Build a Mini Camp",sz=22,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"🧭 探险家任务  Explorer Mission",sz=24,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"Pack · Choose · Build · Safety Rules",sz=14,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"装包 · 选址 · 搭建 · 安全规则",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"野外生存与探险 · Wilderness Survival",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「探险家们, 准备好了吗？」\n• 今天我们不是学生 — 是「小探险家」, 要建一个安全的营地。\n• 4 个任务: 装包 / 搭建 / 选址 / 保护。\n• 完成所有任务可以拿到「探险家徽章」(教师可打印或盖章)。")

# 2. EXPLORER MISSION — 4 tasks
s=ns();bg(s,CREAM);hb(s,"🧭 今天的任务  Today's Mission",PINE)
tb(s,0.4,0.85,9.2,0.45,"🏕️ 我们要去野外露营！",sz=26,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"We're going camping in the wild!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,1.55,9.2,0.32,"👉 为了开心又安全地露营, 我们要先做这 4 件事 ↓",sz=15,b=True,c=BROWN,a=PP_ALIGN.CENTER)
mission_card(s,0.4,1.95,2.25,2.20,1,"装包",     "Pack a backpack", "🎒",SUN)
mission_card(s,2.75,1.95,2.25,2.20,2,"选址",    "Choose location", "📍",SAFCL)
mission_card(s,5.10,1.95,2.25,2.20,3,"搭建",    "Build the camp",  "🏕️",TENTCL)
mission_card(s,7.45,1.95,2.25,2.20,4,"安全规则","Safety Rules",    "🛡️",ALERT)
sentence_frame_bar(s,4.45,
    "为了去露营, 我要 ___ 。",
    "To go camping, I will ___.")
n+=1;pn(s,n)
notes(s,"开场 (1-2 分钟):\n• 先点出今天的大目标: 「我们要去野外露营 camping!」\n• 问: 「想去露营吗? 露营会做什么?」让 2-3 个孩子说。\n• 然后转折: 「但是 — 露营不是说走就走, 我们要先准备!」\n• 介绍 4 件事 (按真实顺序): 装包 → 选址 → 搭建 → 安全规则, 让学生跟读。\n• 重点解释「安全规则」: 不碰野生动物 (像乔治碰臭鼬!), 不一个人玩火, 一直跟着大人。\n• 「完成 1 件打 1 个勾, 4 件都做完, 拿到「探险家徽章」!」\n• 在白板上画 4 格, 每完成一个任务画 ✓。")

# 3. SESSION 1 DIVIDER
s=div("Session 1  上午","📖 Read · 🔍 Discuss · 🏕️ Zones · 🎒 Pack · 📍 Locate",PINE,"📖");n+=1;pn(s,n)

# 5. BEFORE READING — book cover + link + 3 thinking questions
s=ns();bg(s,CREAM);hb(s,"📖 一起读绘本 Read Together · 好奇的乔治去野营",SUN)
# LEFT: book cover image placeholder + link
ib(s,0.4,1.0,3.6,3.2,"📚 绘本封面\nCurious George\nGoes Camping")
link_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.30),Inches(3.6),Inches(0.55))
link_box.fill.solid();link_box.fill.fore_color.rgb=WHITE;link_box.line.color.rgb=SUN;link_box.line.width=Pt(1.5)
tb(s,0.5,4.36,0.5,0.4,"🔗",sz=14)
tb(s,0.95,4.34,3.0,0.3,"https://fliphtml5.com/...",sz=10,b=True,c=PINE)
tb(s,0.95,4.56,3.0,0.25,"Book link · 绘本链接",sz=8,c=GRAY)
# RIGHT: 3 thinking questions while reading
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.25),Inches(1.0),Inches(5.4),Inches(3.85))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=SUN;panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.25),Inches(1.0),Inches(5.4),Inches(0.55))
head.fill.solid();head.fill.fore_color.rgb=SUN;head.line.fill.background()
tb(s,4.4,1.08,5.2,0.4,"🤔 乔治去露营 — 边读边想 Think while reading",sz=13,b=True,c=WHITE)
qs=[("1.","他做了什么？","What did he do?"),
    ("2.","对不对？","Right or wrong?"),
    ("3.","惹了什么麻烦？","What trouble did he cause?")]
y=1.75
for num,cn,en in qs:
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(4.45),Inches(y+0.05),Inches(0.4),Inches(0.4))
    badge.fill.solid();badge.fill.fore_color.rgb=SUN;badge.line.fill.background()
    tb(s,4.45,y+0.1,0.4,0.3,num,sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,4.95,y,4.55,0.4,cn,sz=14,b=True,c=DARK)
    tb(s,4.95,y+0.35,4.55,0.3,en,sz=10,c=GRAY)
    y+=0.95
n+=1;pn(s,n)
notes(s,"读绘本前 (1-2 分钟):\n• 「今天我们一起读一本绘本 — 《好奇的乔治去野营》」\n• 把 3 个问题先读一遍, 让孩子带着问题读 (主动阅读, 不是被动听)。\n• 老师可以投影绘本 (用链接打开 FlipHTML5 在线版), 一页一页和孩子一起看。\n• 阅读中间可以暂停, 问: 「乔治在做什么？这样对吗？」\n• 大约 8-10 分钟读完 23 页。")

# 6a. AFTER READING — 6 scenes for DISCUSSION (no answers yet)
cases=[("🏕️","拉绳, 钻帐篷",  "Tent: pulled ropes",        "❌","乱搭会塌, 绊倒人",     ALERT),
       ("🔥","拿水去灭火",    "Fire: poured water",         "❌","火危险, 要找大人",     ALERT),
       ("🌲","一个人跑进森林","Forest: ran in alone",       "❌","会迷路, 有动物",       ALERT),
       ("🦨","靠近臭鼬",      "Skunk: got too close",       "❌","野生动物 = 危险",      ALERT),
       ("💧","掉水里, 爬树",  "Water/tree: fell, climbed",  "❌","水深, 树高都不安全",   ALERT),
       ("🚨","看到火, 喊大人","Fire spotted: told adults",  "✅","这次做对了!",          PINE)]
s=ns();bg(s,CREAM);hb(s,"🤔 你来判断 You Decide · 乔治做了什么？对不对？",SUN)
tb(s,0.4,0.78,9.2,0.30,"先讨论 — 你觉得对不对? 为什么? Discuss first — right or wrong? Why?",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,scene_cn,scene_en,_mark,_reason,_col) in enumerate(cases):
    row=i//3; cidx=i%3
    x=0.35+cidx*3.15
    y=1.18+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.05),Inches(1.50))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BROWN;sh.line.width=Pt(2)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.08),Inches(y+0.08),Inches(0.42),Inches(0.42))
    badge.fill.solid();badge.fill.fore_color.rgb=BROWN;badge.line.fill.background()
    tb(s,x+0.08,y+0.13,0.42,0.3,str(i+1),sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.55,y+0.08,0.6,0.4,em,sz=22)
    tb(s,x+2.45,y+0.08,0.55,0.45,"❓",sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+0.58,2.90,0.30,scene_cn,sz=12,b=True,c=DARK)
    tb(s,x+0.10,y+0.86,2.90,0.24,scene_en,sz=8,c=GRAY)
    tb(s,x+0.10,y+1.16,2.90,0.30,"对吗? 为什么?",sz=10,b=True,c=BROWN)
sentence_frame_bar(s,4.55,
    "我觉得乔治 ___ , 因为 ___ 。",
    "I think George ___, because ___.")
n+=1;pn(s,n)
notes(s,
"讨论环节 (5-7 分钟) — 先想再说, 不直接给答案:\n"
"• 老师指着每张卡, 让学生先讨论 (同桌 / 小组 / 举手投票)。\n"
"• 1 帐篷: 是在帮忙还是捣乱? 帐篷区怎么做才安全?\n"
"• 2 火: 火区可以一个人去吗? 看到火问题, 自己处理还是叫大人?\n"
"• 3 森林: 露营时可以一个人离开营地吗? 想离开应该怎么办?\n"
"• 4 动物: 看到野生动物, 走近还是站远? 为什么不能喂、追?\n"
"• 5 水/树: 水边为什么危险? 在水边谁应该在旁边?\n"
"• 6 喊大人: 这次乔治做对了吗? 哪里做对了?\n"
"• 让学生用句型: 「我觉得乔治 ___ , 因为 ___ 。」\n"
"• 讨论完, 翻页看答案 ✓\n"
"\n"
"━━━━━━━━━━━━━━━━━━━━━━\n"
"📖 教师参考 — 绘本错误案例分析\n"
"━━━━━━━━━━━━━━━━━━━━━━\n"
"这个绘本很适合当 Day 2 的「错误案例分析」: 乔治想帮忙, 但因为不懂露营规则, 惹出很多麻烦。\n"
"主题问题: 乔治去露营, 他做了什么? 对不对? 惹了什么麻烦?\n"
"\n"
"1️⃣ 他想帮忙搭帐篷\n"
"  情节: 乔治看到大家在搭帐篷, 也想帮忙。他拉绳子、钻进帐篷里玩。\n"
"  问题: 他是真的在帮忙吗?\n"
"  结论: 不太对。\n"
"  原因: 搭帐篷需要按照步骤来, 不能乱拉绳子、乱钻进去, 不然帐篷搭不好, 也可能绊倒别人。\n"
"  课堂问题: 乔治是在帮忙, 还是在捣乱? 帐篷区应该怎么做才安全?\n"
"\n"
"2️⃣ 他去打水, 又想帮忙灭火\n"
"  情节: 乔治看到营地里有水、有火, 就拿水去帮忙。后来他把水倒在火上。\n"
"  问题: 看到火, 能不能自己处理?\n"
"  结论: 不对。\n"
"  原因: 火很危险, 孩子不能自己靠近火、玩火、倒水灭火, 应该告诉大人。\n"
"  课堂问题: 火区可以一个人去吗? 如果看到火有问题, 应该自己处理还是叫大人?\n"
"\n"
"3️⃣ 他一个人跑进森林\n"
"  情节: 乔治听到声音、看到动物, 就跟着跑进森林。\n"
"  问题: 露营时可以一个人离开营地吗?\n"
"  结论: 不对。\n"
"  原因: 森林里容易迷路, 也可能遇到动物、河水、危险地形。\n"
"  课堂问题: 乔治为什么会迷路? 如果你想离开营地, 应该怎么办?\n"
"\n"
"4️⃣ 他靠近动物, 结果被臭鼬喷了\n"
"  情节: 乔治看到动物很好奇, 想靠近看。结果遇到臭鼬, 被喷了一身臭味。\n"
"  问题: 看到野生动物可以靠近吗?\n"
"  结论: 不可以。\n"
"  原因: 野生动物不是宠物, 靠近可能会受伤, 也会打扰动物。\n"
"  课堂问题: 看到动物, 我们应该走近看, 还是站远一点? 为什么不能喂动物、追动物?\n"
"\n"
"5️⃣ 他掉进水里, 又爬到树上\n"
"  情节: 乔治后来靠近水边, 掉进水里, 又爬到树上躲起来。\n"
"  问题: 水边可以自己玩吗? 树上安全吗?\n"
"  结论: 不安全。\n"
"  原因: 水边会滑, 水可能很深; 爬树也可能摔下来。\n"
"  课堂问题: 水边为什么危险? 如果在水边, 谁应该在旁边?\n"
"\n"
"6️⃣ 最后他帮忙发现了火情 ✓\n"
"  情节: 乔治在树上看到远处有烟和火, 提醒大家, 最后大家来灭火。\n"
"  问题: 这次他做对了吗?\n"
"  结论: 这次做对了。\n"
"  原因: 他发现危险后没有自己去处理, 而是让大人知道, 帮助大家发现问题。\n"
"  课堂问题: 看到危险时, 乔治应该自己去救火吗? 正确做法是什么?\n"
"\n"
"━━━━━━━━━━━━━━━━━━━━━━\n"
"💡 总结一句话给学生:\n"
"乔治很热心, 也想帮忙, 但是露营不能只靠好奇心。我们要学会:\n"
"装好背包、分好营地区域、选安全地点、遇到危险找大人。")

# 6b. AFTER READING — 6 cases ANSWERS (reveal)
s=ns();bg(s,CREAM);hb(s,"🔍 答案 The Answers · 乔治做了什么？对不对？",SUN)
tb(s,0.4,0.78,9.2,0.30,"乔治很热心, 想帮忙 — 可是他懂露营规则吗？",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
for i,(em,scene_cn,scene_en,mark,reason_cn,col) in enumerate(cases):
    row=i//3; cidx=i%3
    x=0.35+cidx*3.15
    y=1.18+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.05),Inches(1.50))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=col;sh.line.width=Pt(2)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.08),Inches(y+0.08),Inches(0.42),Inches(0.42))
    badge.fill.solid();badge.fill.fore_color.rgb=col;badge.line.fill.background()
    tb(s,x+0.08,y+0.13,0.42,0.3,str(i+1),sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.55,y+0.08,0.6,0.4,em,sz=22)
    tb(s,x+2.45,y+0.08,0.55,0.45,mark,sz=22,b=True,c=col,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+0.58,2.90,0.30,scene_cn,sz=12,b=True,c=DARK)
    tb(s,x+0.10,y+0.86,2.90,0.24,scene_en,sz=8,c=GRAY)
    tb(s,x+0.10,y+1.14,2.90,0.30,reason_cn,sz=11,b=True,c=col)
sentence_frame_bar(s,4.55,
    "乔治不安全, 因为 ___ 。 应该 ___ 。",
    "George wasn't safe because ___. He should ___.")
n+=1;pn(s,n)
notes(s,"揭晓答案 (3-4 分钟):\n• 一张一张揭开, 看学生猜对了几个。\n• 重点: 5 个 ❌ + 1 个 ✅。\n• 关键提问: 第 6 件事 — 为什么这次乔治是对的? (因为他没有自己处理, 找了大人!)\n• 总结: 乔治很热心, 但热心 ≠ 帮忙 — 要懂规则。\n• 看到危险找大人 — 是英雄, 不是「捣乱」!")

# 7. SYNTHESIS + transition into 4 tasks
s=ns();bg(s,CREAM);hb(s,"🌟 我们学到 What We Learned · 4 个任务",PINE)
# Transition pill (George → rules)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(0.55))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=ALERT;sh.line.width=Pt(2)
tb(s,0.55,0.92,8.9,0.4,"👉 乔治惹麻烦, 因为他不懂规则! 我们来学这 4 件事 →",sz=15,b=True,c=ALERT,a=PP_ALIGN.CENTER)
tb(s,0.55,1.18,8.9,0.28,"George got into trouble because he didn't know the rules — let's learn 4 things →",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tasks=[("🎒","装包",     "Pack a backpack", SUN),
       ("📍","选址",     "Choose location", SAFCL),
       ("🏕️","搭建",     "Build the camp",  TENTCL),
       ("🛡️","安全规则", "Safety Rules",    ALERT)]
for i,(em,cn,en,c) in enumerate(tasks):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.20),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=c;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.85,2.1,0.9,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.85,2.1,0.5,cn,sz=20,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.4,2.1,0.4,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.55,3.85,1.1,0.3,f"任务 {i+1}",c,sz=10)
sentence_frame_bar(s,4.4,
    "我学到 ___ 。下次露营, 我会 ___ 。",
    "I learned ___. Next time camping, I will ___.")
n+=1;pn(s,n)
notes(s,"过渡 (2-3 分钟):\n• 「乔治惹麻烦, 因为他不懂规则! 我们来学规则。」\n• 跟读 4 个任务: 装包 → 选址 → 搭建 → 安全规则。\n• 「先来看「搭建」— 一个好营地有 4 个区!」(下一页)\n• 完成所有任务的学生 → 颁发「探险家徽章」。")

# 8. CAMP STRUCTURE — 5 zones (4 from Shepherd Camp + recreation area)
s=ns();bg(s,CREAM);hb(s,"🏕️ 营地有 5 个区  Camp Has 5 Zones",TENTCL)
tb(s,0.4,0.85,9.2,0.36,"营地像家一样, 有 5 个区 ↓",sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.28,"A camp is like a home — 5 zones ↓",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
zones=[("⛺","帐篷区",     "Tent",      "睡觉休息",   "sleep",     TENTCL),
       ("🔥","用火就餐区","Fire & Eat","烧饭吃饭",   "cook/eat",  FIRECL),
       ("⚽","公共娱乐区","Recreation","玩游戏",     "play/games",PLAYCL),
       ("💧","取水用水区","Water",     "喝水洗碗",   "drink/wash",SAFCL),
       ("🚿","卫生区",     "Sanitation","洗手上厕所","wash/toilet",PINE)]
for i,(em,cn,en,use_cn,use_en,c) in enumerate(zones):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.78),Inches(2.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=c;sh.line.width=Pt(2.5)
    tb(s,x+0.04,1.65,1.70,0.75,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,x+0.04,2.45,1.70,0.40,cn,sz=13,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.04,2.85,1.70,0.28,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.04,3.25,1.70,0.30,use_cn,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.04,3.55,1.70,0.25,use_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.34,3.92,1.10,0.28,f"区 {i+1}",c,sz=9)
sentence_frame_bar(s,4.45,
    "营地有 ___ 个区。这是 ___ 区。",
    "The camp has ___ zones. This is the ___ zone.")
n+=1;pn(s,n)
notes(s,"介绍 5 个区 (3-4 分钟) — 露营营地的基本区域:\n• 「营地像家一样 — 有客厅、厨房、卧室、厕所、操场。露营也一样!」\n• 一个一个介绍, 让学生跟读: 帐篷区 → 用火就餐区 → 公共娱乐区 → 取水用水区 → 卫生区。\n• 用手指数 1-2-3-4-5。\n• 问: 「你睡觉去哪个区? 吃饭去哪个区? 玩游戏去哪个区? 喝水去哪个区? 上厕所去哪个区?」\n• 关键安全点 (5 个区要分开):\n  - 帐篷区: 平整、干净, 清除石块树根\n  - 用火就餐区: 离帐篷远, 注意防火\n  - 公共娱乐区: 大家一起玩, 不能太靠近火和水\n  - 取水用水区: 干净水源, 远离厕所\n  - 卫生区: 离帐篷和水源都要远, 防止污染\n• 下一页: 物品归类活动, 检查理解。\n• (来源参考: Shepherd Camp 露营营地建设四大基本区域 + 公共娱乐区)")

# 9. OBJECT MATCHING — check understanding of 5 zones
s=ns();bg(s,CREAM);hb(s,"🎯 物品归类  Object Match · 这个放哪个区?",SUN)
tb(s,0.4,0.78,9.2,0.30,"看一看 — 这些东西放哪个区?",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.05,9.2,0.24,"Where does each thing go?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 10 items in 2 rows of 5
items=[("🛌","睡袋",  "sleeping bag"),
       ("🔦","手电",  "flashlight"),
       ("🍳","锅",    "pot"),
       ("🍡","棉花糖","marshmallow"),
       ("⚽","球",    "ball"),
       ("🪢","跳绳",  "jump rope"),
       ("💧","水壶",  "water bottle"),
       ("🪣","水桶",  "water bucket"),
       ("🪥","牙刷",  "toothbrush"),
       ("🧼","香皂",  "soap")]
for i,(em,cn,en) in enumerate(items):
    row=i//5; cidx=i%5
    x=0.4+cidx*1.88
    y=1.35+row*0.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(1.78),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BROWN;sh.line.width=Pt(1.5)
    tb(s,x+0.06,y+0.05,0.55,0.55,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.62,y+0.05,1.14,0.28,cn,sz=11,b=True,c=DARK)
    tb(s,x+0.62,y+0.32,1.14,0.25,en,sz=7,c=GRAY)
# 5 zone targets at bottom
zone_targets=[("⛺","帐篷区",  TENTCL),
              ("🔥","用火就餐区",FIRECL),
              ("⚽","公共娱乐区",PLAYCL),
              ("💧","取水用水区",SAFCL),
              ("🚿","卫生区",  PINE)]
for i,(em,cn,c) in enumerate(zone_targets):
    x=0.4+i*1.88
    y=2.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(1.78),Inches(0.95))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,x+0.04,y+0.08,1.70,0.42,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.04,y+0.52,1.70,0.35,cn,sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.10,
    "___ 放在 ___ 区。",
    "The ___ goes in the ___ zone.")
n+=1;pn(s,n)
notes(s,
"检查理解 (5-6 分钟) — 不需要走动:\n"
"• 老师指一个物品, 全班一起说: 「___ 放 ___ 区!」\n"
"• 或者: 老师说物品, 学生指 (用手) 哪个区。\n"
"• 答案 (供参考):\n"
"  - 🛌 睡袋, 🔦 手电 → 帐篷区 (睡觉用)\n"
"  - 🍳 锅, 🍡 棉花糖 → 用火就餐区 (烧、煮、吃)\n"
"  - ⚽ 球, 🪢 跳绳 → 公共娱乐区 (玩游戏)\n"
"  - 💧 水壶, 🪣 水桶 → 取水用水区 (喝水、洗碗)\n"
"  - 🪥 牙刷, 🧼 香皂 → 卫生区 (洗手、刷牙)\n"
"• 难的可以让学生讨论 — 老师追问: 「为什么放在那里? 不放别的区可以吗?」\n"
"• K 句型: 「这是 ___」 + 指。\n"
"• G1-3 句型: 「___ 放在 ___ 区, 因为 ___ 。」")

# ===== 5 ZONE ANSWER SLIDES (direct answer with rich CN+EN bullets) =====
# Each slide auto-uses zone_*.png from pics/ if it exists (drop AI-generated images there)
_PICS = "/Users/Huan/0 projects/summercourse/Chinese/野外生存与探险wilderness_pbl/pics"

# Zone 1: 帐篷区
n+=1;s=answer_panels_slide(
    "⛺ 帐篷区 · 💡 露营小贴士 Pro Tips",
    TENTCL,
    [
        {"q":"帐篷应该搭在什么样的地方?","mark":"✅","lines":[
            "帐篷应该搭在平坦、干燥、安全的地方 Flat, dry, safe ground",
            "离河边不要太近 Not too close to the river",
            "没有大石头和树枝 No big rocks or branches",
        ]},
        {"q":"下雨、起风时, 要注意什么?","mark":"✅","lines":[
            "下雨天小心水会流进帐篷 Rain may flow into tent",
            "地面太湿会积水 Wet ground may flood",
            "风大要把帐篷固定好 Strong wind — secure it well",
        ]},
        {"q":"可以搭在大树正下方吗?","mark":"❌","lines":[
            "树枝可能掉下来 Branches may fall",
            "打雷时不安全 Not safe during storms",
            "下雨时会一直滴水 Rain keeps dripping",
        ]},
    ],
    img_label="📷 平整干净的好地点",
    img_path=os.path.join(_PICS,"zone_tent.png"))
pn(s,n)

# Zone 2: 用火就餐区
n+=1;s=answer_panels_slide(
    "🔥 用火就餐区 · 💡 露营小贴士 Pro Tips",
    FIRECL,
    [
        {"q":"火堆应该在哪里, 怎么搭?","mark":"✅","lines":[
            "火堆离帐篷至少 4 米远 Fire pit 4m+ from tents",
            "用石头围起来 + 准备水或沙 Stones around + water/sand ready",
            "大人陪同才能用火 Adult must be there",
        ]},
        {"q":"吃完饭、走之前要做什么?","mark":"✅","lines":[
            "食物垃圾装袋带走 Pack food trash in a bag",
            "锅碗洗干净, 不留食物味 Wash cookware",
            "用水或沙完全压灭火 Fully extinguish",
        ]},
        {"q":"哪些事情绝对不能做?","mark":"❌","lines":[
            "不能一个人离开火堆 Never leave fire alone",
            "食物垃圾不能扔地上 — 熊会来 No scraps on ground (bears!)",
            "不烧塑料 — 有毒 + 会爆炸 No plastic in fire",
        ]},
    ],
    img_label="📷 用石头围起来的火堆",
    img_path=os.path.join(_PICS,"zone_fire.png"))
pn(s,n)

# Zone 3: 公共娱乐区
n+=1;s=answer_panels_slide(
    "⚽ 公共娱乐区 · 💡 露营小贴士 Pro Tips",
    PLAYCL,
    [
        {"q":"玩游戏在哪里玩才安全?","mark":"✅","lines":[
            "玩游戏要离火堆和帐篷远 Far from fire and tents",
            "在大人能看到的地方玩 Where adults can see",
            "走小路, 不踩花草 Stay on the trail",
        ]},
        {"q":"玩耍时还要注意什么?","mark":"✅","lines":[
            "天黑前回到营地 Back to camp before dark",
            "声音小一点, 不打扰别人 Keep voice down",
            "和小伙伴一起玩 Buddy up — never alone",
        ]},
        {"q":"哪些事情不能做?","mark":"❌","lines":[
            "天黑不乱跑 — 看不到地面 No running after dark",
            "不爬树, 不摘花 No climbing / picking",
            "不打扰野生动物 Don't disturb wildlife",
        ]},
    ],
    img_label="📷 安全玩耍的地方",
    img_path=os.path.join(_PICS,"zone_recreation.png"))
pn(s,n)

# Zone 4: 取水用水区
n+=1;s=answer_panels_slide(
    "💧 取水用水区 · 💡 露营小贴士 Pro Tips",
    SAFCL,
    [
        {"q":"喝水之前要怎么做?","mark":"✅","lines":[
            "水从干净的水源取 Take from clean source",
            "喝水前要烧开或过滤 Boil OR filter first",
            "每人都有自己的水壶 Each person has own bottle",
        ]},
        {"q":"洗碗、用水要注意什么?","mark":"✅","lines":[
            "洗碗在水区, 不在帐篷旁 Wash dishes here, not at tent",
            "水区要远离厕所 Far from the toilet zone",
            "用桶装脏水 Use bucket for wastewater",
        ]},
        {"q":"哪些事情绝对不能做?","mark":"❌","lines":[
            "不能直接喝生水 — 有寄生虫 Never drink raw water!",
            "不在水边打闹 — 容易滑倒 No play at water",
            "不在河里用肥皂 — 污染水源 No soap in the river",
        ]},
    ],
    img_label="📷 滤水器/烧水的水壶",
    img_path=os.path.join(_PICS,"zone_water.png"))
pn(s,n)

# Zone 5: 卫生区
n+=1;s=answer_panels_slide(
    "🚿 卫生区 · 💡 露营小贴士 Pro Tips",
    PINE,
    [
        {"q":"厕所应该在哪里? 怎么洗手?","mark":"✅","lines":[
            "厕所离帐篷和水源都 60 米以上 Toilet 60m+ from tents/water",
            "用免洗洗手液 + 湿纸巾 Hand sanitizer + wet wipes",
            "吃饭前洗手 Wash hands before meals",
        ]},
        {"q":"垃圾、人体垃圾怎么处理?","mark":"✅","lines":[
            "上厕所挖小坑 (约 15 厘米深) 埋好 Dig hole ~6 inches and bury",
            "湿纸巾、垃圾全部装袋带走 Pack out ALL wipes and trash",
            "营地保持干净 — 像没人来过 Leave No Trace",
        ]},
        {"q":"哪些事情绝对不能做?","mark":"❌","lines":[
            "不留下任何垃圾 No trace left behind",
            "不在水源附近上厕所 Never near water source",
            "不在帐篷里穿外面的鞋 No outdoor shoes inside tent",
        ]},
    ],
    img_label="📷 露营厕所/Leave No Trace",
    img_path=os.path.join(_PICS,"zone_sanitation.png"))
pn(s,n)

# ===== 5 "Can we camp HERE?" SCENARIOS — direct-answer format =====

# 1. 大石头 / 悬崖
n+=1;s=answer_panels_slide(
    "🔍 这里能搭营吗？  Can we camp HERE? — 大石头",
    ALERT,
    [
        {"q":"为什么这里不安全?","mark":"❌","lines":[
            "上面有悬崖, 石头可能掉下来 Rocks may fall",
            "风吹、雨水会让石头松动 Wind/rain loosen rocks",
            "帐篷会被砸破, 人会受伤 Tent + people get hurt",
        ]},
        {"q":"还有什么麻烦?","mark":"❌","lines":[
            "晚上看不见石头, 容易摔倒 Hard to see at night",
            "石头会反射热量 / 寒气 Rocks reflect heat/cold",
        ]},
        {"q":"那应该选什么样的地方?","mark":"✅","lines":[
            "应该选远离悬崖的平地 Flat ground, far from cliffs",
            "上面没有掉落物 No falling hazards overhead",
            "至少离大石头 30 米 At least 30m from big rocks",
        ]},
    ],
    img_label="📷 camp_rocks_scenario.png",
    img_path=os.path.join(_PICS,"camp_rocks_scenario.png"),
    subtitle="揭晓答案 — 这里不安全! Look at the picture — it's NOT safe.")
pn(s,n)

# 2. 草地 / 虫子
n+=1;s=answer_panels_slide(
    "🔍 这里能搭营吗？(2)  Can we camp HERE? — 草地",
    ALERT,
    [
        {"q":"草丛里有什么不好?","mark":"❌","lines":[
            "草丛里有虫子 (蚂蚁、蚊子、蜘蛛) Bugs in the grass",
            "虫子会爬进帐篷, 咬人 Bugs may crawl into tent",
            "高草看不到地面石头 Tall grass hides hazards",
        ]},
        {"q":"下雨、早上又会怎样?","mark":"❌","lines":[
            "早上有露水, 帐篷底潮湿 Dew makes tent floor wet",
            "下雨后泥泞难走 Muddy after rain",
        ]},
        {"q":"那应该选什么样的地方?","mark":"✅","lines":[
            "应该选短草 / 干泥土 / 沙地 Short grass / dirt / sand",
            "搭帐篷前先清扫地面 Sweep area first",
            "检查有没有蚂蚁洞 Check for ant nests",
        ]},
    ],
    img_label="📷 camp grass scenario.png",
    img_path=os.path.join(_PICS,"camp grass scenario.png"),
    subtitle="揭晓答案 — 这里不舒服! Look at the picture — it's NOT comfortable.")
pn(s,n)

# 3. 大树正下方
n+=1;s=answer_panels_slide(
    "🔍 这里能搭营吗？(3)  Can we camp HERE? — 大树下",
    ALERT,
    [
        {"q":"为什么这里不安全?","mark":"❌","lines":[
            "枯枝可能掉下来 Dead branches may fall",
            "打雷时大树有危险 — 不要躲下面 Lightning danger",
            "下雨时树叶会一直滴水 Rain keeps dripping",
        ]},
        {"q":"还有什么不舒服?","mark":"❌","lines":[
            "树根让地面不平 Tree roots make ground uneven",
            "动物会掉下东西 Animals drop things",
        ]},
        {"q":"那应该选什么样的地方?","mark":"✅","lines":[
            "应该在开阔地搭营 (大树附近不在下方) Near trees, not under",
            "用大树挡风, 但帐篷在旁边 Tree as windbreak, beside it",
            "看头顶是否有枯枝 Check overhead for dead branches",
        ]},
    ],
    img_label="📷 camp tree scenario.png",
    img_path=os.path.join(_PICS,"camp tree scenario.png"),
    subtitle="揭晓答案 — 这里不安全! Look at the picture — it's NOT safe.")
pn(s,n)

# 4. 海边 / 水边
n+=1;s=answer_panels_slide(
    "🔍 这里能搭营吗？(4)  Can we camp HERE? — 水边",
    ALERT,
    [
        {"q":"为什么这里不安全?","mark":"❌","lines":[
            "涨潮、大浪会冲走帐篷 Tide/waves can wash tent away",
            "半夜涨潮人睡着, 很危险 Tide at night = danger",
            "河水暴涨速度很快 Flash floods rise fast",
        ]},
        {"q":"还有什么麻烦?","mark":"❌","lines":[
            "沙地下雨会塌陷 Wet sand collapses",
            "水边蚊子多 Many mosquitoes near water",
        ]},
        {"q":"那应该选什么样的地方?","mark":"✅","lines":[
            "应该离水至少 30 米 At least 30m from water",
            "选高一点的地方 Choose higher ground",
            "看看「最高水位线」标志 Look for high-water marks",
        ]},
    ],
    img_label="📷 camp water scenario.png",
    img_path=os.path.join(_PICS,"camp water scenario.png"),
    subtitle="揭晓答案 — 这里不安全! Look at the picture — it's NOT safe.")
pn(s,n)

# 5. 低洼地
n+=1;s=answer_panels_slide(
    "🔍 这里能搭营吗？(5)  Can we camp HERE? — 低洼地",
    ALERT,
    [
        {"q":"下雨时会发生什么?","mark":"❌","lines":[
            "雨水从四面流进来 Rain flows in from all sides",
            "帐篷底泡水, 睡袋会湿 Tent floor floods",
            "湿了会冷, 容易生病 Wet = cold = sick",
        ]},
        {"q":"还有什么不舒服?","mark":"❌","lines":[
            "潮湿地方蚊子很多 Damp = mosquitoes",
            "地面软, 帐篷桩固定不牢 Soft ground — stakes loose",
        ]},
        {"q":"那应该选什么样的地方?","mark":"✅","lines":[
            "应该选高地 high ground Higher, flat ground",
            "稍微有斜坡可以排水 Slight slope helps drainage",
            "土要干、要硬 Soil should be dry and firm",
        ]},
    ],
    img_label="📷 low or wet place scenario.png",
    img_path=os.path.join(_PICS,"low or wet place scenario.png"),
    subtitle="揭晓答案 — 这里不舒服! Look at the picture — it's NOT good.")
pn(s,n)

# 16. BACKPACK — INQUIRY: brainstorm before reveal
s=ns();bg(s,CREAM);hb(s,"🤔 露营要带什么？ What to Pack?",SUN)
tb(s,0.4,0.85,9.2,0.4,"如果你要去露营, 你会带什么？",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.32,"If you were going camping, what would you pack?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# Big backpack visual on the left
img_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.85),Inches(4.30),Inches(2.85))
img_box.fill.solid();img_box.fill.fore_color.rgb=WARM;img_box.line.color.rgb=SUN;img_box.line.width=Pt(2.5)
tb(s,0.4,2.40,4.30,1.0,"🎒",sz=110,a=PP_ALIGN.CENTER)
tb(s,0.4,3.65,4.30,0.4,"我的背包 My Backpack",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,0.4,4.05,4.30,0.4,"What goes inside?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Right: prompts
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.85),Inches(4.85),Inches(2.85))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=SUN;panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.85),Inches(4.85),Inches(0.50))
head.fill.solid();head.fill.fore_color.rgb=SUN;head.line.fill.background()
tb(s,5.0,1.92,4.6,0.4,"💭 想一想 — Brainstorm",sz=14,b=True,c=WHITE)
tf=tb(s,5.0,2.50,4.7,0.35,"❓ 你会带什么吃的?",sz=13,b=True,c=DARK)
ap(tf," ",sz=7);ap(tf,"   What food will you bring?",sz=9,c=GRAY);ap(tf," ",sz=8)
ap(tf,"❓ 你会带什么穿的?",sz=13,b=True,c=DARK)
ap(tf,"   What clothes?",sz=9,c=GRAY);ap(tf," ",sz=8)
ap(tf,"❓ 还要带什么 (工具)?",sz=13,b=True,c=DARK)
ap(tf,"   What tools / other things?",sz=9,c=GRAY)
sentence_frame_bar(s,4.80,
    "我会带 ___ 、 ___ 和 ___ 。",
    "I will bring ___, ___, and ___.")
n+=1;pn(s,n)
notes(s,"开放式提问 (3-5 分钟) — 学生先想自己列清单!\n• 「想一想 — 你要去露营 3 天, 你会带什么?」\n• 让学生:\n  - K: 用手指数, 说 1-2 样: 「我带 ___」\n  - G1-3: 列 3-5 样, 用「我会带 ___ 、 ___ 和 ___ 」\n• 老师在白板上记录学生说的所有东西 (不评价对错)。\n• 引导分类: 吃的 / 穿的 / 工具 (有用的)\n• 5 分钟后翻页 — 看老师推荐的 9 件东西, 对比学生的清单。\n• 关键: 不要给答案! 让学生先想, 才会记住。")

# 17. BACKPACK — reveal: 9 essentials
s=ns();bg(s,CREAM);hb(s,"🎒 装什么进背包？  What goes in the backpack?",SUN)
items=[("💧","水","Water"),("🍎","食物","Food"),("🔦","手电筒","Flashlight"),
       ("🧥","外套","Jacket"),("🆘","急救包","First aid"),("🗺️","地图","Map"),
       ("📢","哨子","Whistle"),("🧴","防晒","Sunscreen"),("🧦","袜子","Socks")]
for i,(em,cn,en) in enumerate(items):
    col=i%5;row=i//5
    x=0.4+col*1.92;y=1.0+row*1.7
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(1.78),Inches(1.55))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SUN;sh.line.width=Pt(2)
    tb(s,x+0.05,y+0.05,1.7,0.65,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.78,1.7,0.4,cn,sz=15,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.18,1.7,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.5,
    "我带 ___ 和 ___ 。我不带 ___ 。",
    "I bring ___ and ___. I don't bring ___.")
n+=1;pn(s,n)
notes(s,"先看 9 个东西, 让学生熟悉。下一页: 情况变了 — 不同地方要带不同东西。")

# 19-21. SITUATION CHANGE — Q (think first) → A (reveal) per situation
def situation_q_slide(em,situation_cn,situation_en,color,prompts):
    """Q slide — situation banner + thinking prompts, NO items shown."""
    s=ns();bg(s,CREAM);hb(s,f"{em} 情况变了!  Situation Change!",color)
    big=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(1.4))
    big.fill.solid();big.fill.fore_color.rgb=color;big.line.fill.background()
    tb(s,0.6,1.05,8.8,0.6,situation_cn,sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.6,1.65,8.8,0.4,situation_en,sz=14,c=WARM,a=PP_ALIGN.CENTER)
    tb(s,0.6,2.0,8.8,0.3,"🤔 想一想 — 你要带什么？  Think — what to pack?",sz=14,b=True,c=WARM,a=PP_ALIGN.CENTER)
    # Thinking prompt cards (no answers)
    for i,(p_cn,p_en) in enumerate(prompts):
        x=0.6+i*3.05
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.55),Inches(2.85),Inches(1.55))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
        tb(s,x+0.05,2.65,2.75,0.55,"❓",sz=36,b=True,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.25,2.75,0.40,p_cn,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.70,2.75,0.32,p_en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    sentence_frame_bar(s,4.25,
        "在 ___ , 我会带 ___ 。",
        "In ___, I will bring ___.")
    return s

def situation_a_slide(em,situation_cn,situation_en,color,must_have):
    """A slide — same banner + revealed items."""
    s=ns();bg(s,CREAM);hb(s,f"{em} 答案揭晓!  The Answer!",color)
    big=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(1.4))
    big.fill.solid();big.fill.fore_color.rgb=color;big.line.fill.background()
    tb(s,0.6,1.05,8.8,0.6,situation_cn,sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.6,1.65,8.8,0.4,situation_en,sz=14,c=WARM,a=PP_ALIGN.CENTER)
    tb(s,0.6,2.0,8.8,0.3,"💡 必带的东西 — Must-have items",sz=14,b=True,c=WARM,a=PP_ALIGN.CENTER)
    for i,(emi,cn,en) in enumerate(must_have):
        x=0.6+i*3.05
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.55),Inches(2.85),Inches(1.55))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
        tb(s,x+0.05,2.6,2.75,0.7,emi,sz=44,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.4,2.75,0.4,cn,sz=18,b=True,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.85,2.75,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    sentence_frame_bar(s,4.25,
        "在 ___ 我带 ___ 。",
        "In ___, I bring ___.")
    return s

DESERT_CL = RGBColor(0xD4,0xA5,0x74)
COLD_CL   = RGBColor(0x54,0x6E,0x7A)

# Desert: Q + A
s=situation_q_slide("🏜️","沙漠 — 太阳很大, 没有水","Desert — strong sun, no water",DESERT_CL,
    [("喝的东西?","Drink?"),("挡太阳的?","Sun protection?"),("怕晒怎么办?","Skin?")])
n+=1;pn(s,n)
notes(s,"先讨论 (2-3 分钟):\n• 「沙漠超热, 没水 — 你怎么办?」\n• 不给答案! 让学生想 3 样东西。\n• 引导提示: 「喝什么? 戴什么? 涂什么?」\n• 让 2-3 个学生说: 「在沙漠我会带 ___」")

s=situation_a_slide("🏜️","沙漠 — 太阳很大, 没有水","Desert — strong sun, no water",DESERT_CL,
    [("💧","多带水","Extra water"),("🧴","防晒","Sunscreen"),("🧢","帽子","Hat")])
n+=1;pn(s,n)
notes(s,"揭晓: 沙漠 — 重点是水! 1 天每人需要 4 升水。皮肤防晒。中午躲太阳。")

# Forest + rain: Q + A
s=situation_q_slide("🌧️","森林 + 下雨","Forest + raining",SAFCL,
    [("淋雨怎么办?","If you get wet?"),("脚湿怎么办?","Wet feet?"),("生火湿了?","Wet matches?")])
n+=1;pn(s,n)
notes(s,"先讨论 (2-3 分钟):\n• 「下雨, 树都湿了 — 你怎么办?」\n• 不给答案! 让学生猜 3 样东西。\n• 引导: 「身上 / 脚上 / 火?」")

s=situation_a_slide("🌧️","森林 + 下雨","Forest + raining",SAFCL,
    [("🧥","雨衣","Raincoat"),("👢","胶鞋","Boots"),("🔥","防水火柴","Waterproof matches")])
n+=1;pn(s,n)
notes(s,"揭晓: 下雨 — 鞋子和外衣都要防水。湿衣服 = 失温危险。")

# Cold mountain: Q + A
s=situation_q_slide("🏔️","高山 — 很冷","High mountain — very cold",COLD_CL,
    [("手怎么办?","Hands?"),("脖子冷?","Neck cold?"),("冷想吃什么?","Cold = need?")])
n+=1;pn(s,n)
notes(s,"先讨论 (2-3 分钟):\n• 「山上很冷, 你怎么办?」\n• 不给答案! 让学生猜 3 样东西。\n• 引导: 「手上戴什么? 脖子上呢? 吃什么暖和?」")

s=situation_a_slide("🏔️","高山 — 很冷","High mountain — very cold",COLD_CL,
    [("🧤","手套","Gloves"),("🧣","围巾","Scarf"),("🍫","高能量食物","High-energy food")])
n+=1;pn(s,n)
notes(s,"揭晓: 冷 — 多层衣服比一件厚衣服好 (layering)。手套和围巾保护手指和脖子。高能量食物给身体热量。")

# 23-26. TENT LOCATION A/B (4 rounds)
s=ab_slide("选址 1: 在哪里搭帐篷？","Location 1: where to pitch the tent?",
    "你想在 A 还是 B 搭帐篷？","Where would you pitch the tent — A or B?",
    "🟢","平地","平的, 没有大石头",
    "⛰️","斜坡","斜的, 一边高一边低",
    "A 平地", "睡在斜坡上会滚下去, 帐篷也站不稳。")
n+=1;pn(s,n)

s=ab_slide("选址 2: 离河多远？","Location 2: how far from the river?",
    "在 A 还是 B 搭帐篷？","Pitch the tent at A or B?",
    "💧","河边 (3 米)","就在河旁边",
    "🌳","远离河 (30 米)","远一点的高地",
    "B 远离河", "下大雨河水会上来 (洪水), 河边 3 米的帐篷会被淹。")
n+=1;pn(s,n)

s=ab_slide("选址 3: 在大树下还是空地?","Location 3: under tree or open?",
    "选 A 还是 B?","Choose A or B?",
    "🌲","大树下","有阴凉但有枯枝",
    "🌿","空地","没遮挡但安全",
    "B 空地", "大树底下的枯枝叫「widowmaker」, 风一吹会掉下来。下雨打雷时, 大树会引雷。")
n+=1;pn(s,n)

s=ab_slide("选址 4: 离火堆多远?","Location 4: how far from fire?",
    "选 A 还是 B?","Choose A or B?",
    "🏕️🔥","靠近火 (1 米)","暖和, 但火星会飞过来",
    "🏕️➡️🔥","离火 4 米","暖和靠走过去, 不会着火",
    "B 离火 4 米", "火星可能飞 1-2 米, 把帐篷烧着。安全距离 = 4 米以上。")
n+=1;pn(s,n)

# 27. REASONING SUMMARY
s=ns();bg(s,CREAM);hb(s,"🧠 你说得对!  You Got It!",PINE)
tb(s,0.4,0.9,9.2,0.4,"4 个选址规则:  4 location rules:",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
rules=[("🟢 1.","平地 — 不要斜坡","Flat — no slope"),
       ("💧 2.","离河 30 米 — 不要被洪水冲","30m from river — avoid floods"),
       ("🌿 3.","空地 — 不在大树下","Open — not under big trees"),
       ("🔥 4.","离火堆 4 米以上","4m+ from the fire")]
for i,(em,cn,en) in enumerate(rules):
    y=1.5+i*0.78
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(2)
    tb(s,0.6,y+0.1,1.4,0.5,em,sz=20,b=True,c=PINE)
    tb(s,2.0,y+0.05,5.4,0.35,cn,sz=15,b=True,c=DARK)
    tb(s,2.0,y+0.4,5.4,0.3,en,sz=11,c=GRAY)
n+=1;pn(s,n)

# 28. SESSION 2.5 — Video interactions divider
s=div("🎬 看视频学探险","Watch & Learn  ·  Camping setup · Safe campsite · Backpack",SAFCL,"📺");n+=1;pn(s,n)

# 29. VIDEO 1 — Camping setup
s=video_slide("How to Set Up a Tent (Kids)","怎么搭帐篷",
    "👂 听 / 看:\n1. 搭帐篷的第一步是什么？\n2. 一共有几步？\n3. 需要几个人?",
    "🎯 看完后:\n• 演一演: 跟同学一起「假装搭帐篷」\n• 数一数: 共 ___ 步\n• 说一说: 第一步是 ___",
    bgc=PINE);n+=1;pn(s,n)
notes(s,"视频建议: YouTube 搜「How to set up a tent for kids」(选 1-3 分钟版本)。\n关键步骤: 1) 选地方 2) 铺地布 3) 立支架 4) 拉绳子。")

# 30. VIDEO 2 — Safe vs unsafe campsite
s=video_slide("Safe vs Unsafe Campsite","安全 vs 不安全的营地",
    "👂 听 / 看:\n1. 这个营地哪里安全？\n2. 哪里不安全？\n3. 你能找到 3 个错？",
    "🎯 看完后:\n• 圈一圈: 在 worksheet 上画 ✅ 和 ❌\n• 说一说: 「___ 不安全, 因为 ___」\n• 决定: 让乔治再来吗？",
    bgc=ALERT);n+=1;pn(s,n)
notes(s,"视频建议: YouTube 搜「camping safety for kids」或 Smokey Bear 的视频。或老师自己拍 30 秒手机视频。")

# 31. VIDEO 3 — Packing backpack
s=video_slide("Packing a Backpack","怎么装背包",
    "👂 听 / 看:\n1. 哪些东西先放进去？\n2. 哪些放在最上面？\n3. 重的放哪里？",
    "🎯 看完后:\n• 排顺序: 1)___ 2)___ 3)___\n• 说一说: 重的东西放 ___\n• 比一比: 跟你装的一样吗?",
    bgc=SUN);n+=1;pn(s,n)
notes(s,"原则: 重的放中间靠背 (重心稳)。睡袋放下面。常用的 (水/急救) 放最上。")

# 32. SESSION 2 DIVIDER
s=div("Session 2  下午","🔄 复习 + 语言目标 (我会认/我会写)",PINE,"📖");n+=1;pn(s,n)

# 33. REVIEW — Baamboozle game intro
s=ns();bg(s,CREAM);hb(s,"🎮 复习  Review · Baamboozle Game",SUN)
tb(s,0.4,0.85,9.2,0.45,"今天学了好多 — 来玩 Baamboozle 复习!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.30,"We learned a lot today — let's review with Baamboozle!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 4 review topic cards
topics=[("🏕️","5 个营地区","5 Zones",         SAFCL),
        ("🎒","背包要带什么","Backpack items",  SUN),
        ("📍","帐篷搭哪里","Tent location",    TENTCL),
        ("🐵","乔治的故事","George's story",   ALERT)]
for i,(em,cn,en,c) in enumerate(topics):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(2.20),Inches(1.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.80,2.10,0.65,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.50,2.10,0.40,cn,sz=14,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.92,2.10,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.40,3.25,1.40,0.25,"会答!",c,sz=10)
# Game info banner
gb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(0.95))
gb.fill.solid();gb.fill.fore_color.rgb=WARM
gb.line.color.rgb=SUN;gb.line.width=Pt(2)
tb(s,0.55,3.95,9.0,0.40,"🎯 怎么玩 How to play:",sz=14,b=True,c=SUN)
tb(s,0.55,4.30,9.0,0.30,"分组 → 老师点开 Baamboozle → 学生答对得分 → 最高分赢!",sz=12,b=True,c=DARK)
tb(s,0.55,4.55,9.0,0.25,"Teams → teacher opens Baamboozle → answer to score → highest wins!",sz=10,c=GRAY)
n+=1;pn(s,n)
notes(s,"复习游戏 (10-15 分钟):\n• 老师提前在 Baamboozle 上 import 「day2_review_baamboozle.csv」 文件 (在同一文件夹).\n• 20 个问题, 涵盖 5 营区 / 背包 / 帐篷地点 / 乔治故事。\n• 分 2-4 组玩。每组答对加分。\n• 答错的问题, 全班一起复习正确答案。\n• Baamboozle 网址: https://www.baamboozle.com/")

# 34. VOCAB OVERVIEW
s=ns();bg(s,CREAM);hb(s,"📚 今天的字  Today's Words",PINE)
tb(s,0.4,0.9,9.2,0.4,"我会认 5 个词  ·  我会写 3 个字",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
read_words=[("营","yíng","camp"),("帐篷","zhàng peng","tent"),("火","huǒ","fire"),
            ("包","bāo","bag"),("安全","ān quán","safe")]
write_chars=[("火","huǒ","fire"),("包","bāo","bag"),("水","shuǐ","water")]
# Read row
tb(s,0.4,1.55,9.2,0.4,"👀 我会认  I Can Read",sz=16,b=True,c=SUN)
for i,(w,py,en) in enumerate(read_words):
    x=0.4+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(1.78),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2)
    tb(s,x+0.05,2.05,1.7,0.55,w,sz=28,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,1.7,0.3,py,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.92,1.7,0.3,en,sz=10,c=DARK,a=PP_ALIGN.CENTER)
# Write row
tb(s,0.4,3.4,9.2,0.4,"✍️ 我会写  I Can Write",sz=16,b=True,c=PINE)
for i,(w,py,en) in enumerate(write_chars):
    x=2.0+i*2.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.85),Inches(1.9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(2)
    tb(s,x+0.05,3.92,1.8,0.55,w,sz=32,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.55,1.8,0.3,py,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.82,1.8,0.3,en,sz=10,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# 35-39. WORD CARDS — 我会认
read_data=[
    ("营","yíng","camp",          "我们去搭一个营地。",       "📷 营地 / Tent + fire camp"),
    ("帐篷","zhàng peng","tent",  "我的帐篷是绿色的。",       "📷 绿色帐篷"),
    ("火","huǒ","fire",            "火很热, 不要靠近。",       "📷 营火 (用石头围)"),
    ("包","bāo","bag",             "我的背包里有水和食物。",   "📷 背包 + 物品"),
    ("安全","ān quán","safe",      "在大人旁边最安全。",       "📷 大人 + 孩子在营地"),
]
for w,py,en,sent,img in read_data:
    s=word_card_read(w,py,en,sent,img,color=SUN);n+=1;pn(s,n)

# 40-42. WORD CARDS — 我会写
write_data=[
    ("火","huǒ","fire",  "4 笔: 点 - 撇 - 撇 - 捺\n像火苗的样子! 🔥"),
    ("包","bāo","bag",   "5 笔: 撇 - 横折钩 - 横折 - 横 - 竖弯钩\n外面像「勹」, 里面是「巳」"),
    ("水","shuǐ","water","4 笔: 竖钩 - 横撇 - 撇 - 捺\n中间一竖, 两边像水滴 💧"),
]
for w,py,en,strokes in write_data:
    s=word_card_write(w,py,en,strokes,color=PINE);n+=1;pn(s,n)

# 43. SESSION 3 DIVIDER
s=div("Session 3  下午","🎒 写 Booklet + 2 个项目  Booklet + Camp Model + Pack Challenge",BROWN,"🛠️");n+=1;pn(s,n)

# 44. Day 2 Booklet
s=ns();bg(s,CREAM);hb(s,"📓 完成 Day 2 练习册  Day 2 Booklet",BROWN)
tb(s,0.4,0.85,9.2,0.35,"老师带着一起做  Teacher leads — do it together",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
ib(s,0.4,1.25,9.2,4.0,"📷 Booklet 截图")
n+=1;pn(s,n)
notes(s,"老师带着学生一页一页完成 Day 2 练习册。可以用「先看 → 一起读 → 自己做 → 同桌检查」。")

# 45. Projects overview — 2 cards side by side
s=ns();bg(s,CREAM);hb(s,"🛠️ 动手时间！  Hands-On Time — 选一个项目",BROWN)
projects=[
    ("PROJECT 1","🏕️ 最安全的露营地","Camp Model — Design the safest camp",
     "适合高 level · 手工有难度",TENTCL),
    ("PROJECT 2","🎒 露营背包挑战","Pack Your Backpack — Pack the right items",
     "适合低 level · 简单上手",SUN),
]
for i,(lbl,nm,en,lvl,cl) in enumerate(projects):
    x=0.4+i*4.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(4.4),Inches(4.15))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,4.2,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.4,4.2,0.55,nm,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.95,4.2,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+1.2,2.35,2.0,0.35,lvl.split(" · ")[0],cl,sz=11)
    ib(s,x+0.2,2.85,4.0,1.5,"📷 示范")
    tf=tb(s,x+0.15,4.45,4.1,0.35,lvl.split(" · ")[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    ap(tf,lvl.split(" · ")[1] if " · " in lvl else "",sz=12,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"两个项目根据班级 level 二选一。Project 1 (Camp Model) 需要更多手工材料和时间, 适合 G1-3。Project 2 (Pack Challenge) 用图卡, 简单, 适合 K-G1。")

# 46. Project 1 — Camp Model (high-level)
s=ns();bg(s,CREAM);hb(s,"🏕️ Project 1: 最安全的露营地  Camp Model",TENTCL)
# Left: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=TENTCL;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.3,"📋 纸板 (底板)  Cardboard base",sz=13,c=DARK)
ap(tf,"📦 小盒子 (帐篷)  Small box (tent)",sz=13,c=DARK)
ap(tf,"📄 纸 + 彩笔  Paper + markers",sz=13,c=DARK)
ap(tf,"📎 胶带  Tape",sz=13,c=DARK)
# Right: steps
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.3,"1️⃣ 选环境: 🌊 海边 or 🌳 森林?  Choose env",sz=13,c=DARK)
ap(tf2,"2️⃣ 摆出 4 个区域: ⛺ 帐篷 / 🔥 生火区 (远离帐篷!) / 🍱 食物区 / 🟢 安全区",sz=12,c=DARK)
ap(tf2,"3️⃣ 用纸/盒子做出来  Build it",sz=13,c=DARK)
ap(tf2,"4️⃣ 跟同学讲一讲  Explain to a friend",sz=13,c=DARK)
# Bottom: sentence frames
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.9),Inches(9.4),Inches(1.25))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=TENTCL;sh3.line.width=Pt(2)
tb(s,0.5,4.0,9,0.35,"🗣️ 展示句型:",sz=14,b=True,c=TENTCL)
tb(s,0.5,4.4,9,0.35,"· 这是我的营地。",sz=14,c=DARK)
tb(s,0.5,4.7,9,0.35,"· 我选了 ___。 (海边 / 森林)",sz=14,c=DARK)
tb(s,5.0,4.4,4.7,0.35,"· 我把 ___ 放在 ___, 因为 ___。",sz=14,c=DARK)
n+=1;pn(s,n)
notes(s,"高 level 项目: 用纸盒、纸板、彩笔做一个 3D 营地模型。重点不在精美, 而在 4 区分开 + 能讲出原因。")

# 47. Project 2 — Backpack Challenge (low-level)
s=ns();bg(s,CREAM);hb(s,"🎒 Project 2: 露营背包挑战  Pack Your Backpack",SUN)
# Left: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(3.0),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=SUN;sh.line.fill.background()
tb(s,0.4,0.98,2.8,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,3.0,2.3,"🃏 物品图卡  Item picture cards (or worksheet)",sz=12,c=DARK)
ap(tf,"✏️ 铅笔 / Pencil",sz=12,c=DARK)
ap(tf,"✅ ❌ 贴纸 (or stickers)",sz=12,c=DARK)
# Middle: items
sh_ex=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.45),Inches(0.95),Inches(3.0),Inches(0.4))
sh_ex.fill.solid();sh_ex.fill.fore_color.rgb=BROWN;sh_ex.line.fill.background()
tb(s,3.55,0.98,2.8,0.35,"💡 物品  Items",sz=14,b=True,c=WHITE)
tf_ex=tb(s,3.55,1.45,3.0,2.3,"✔️ 有用: 帐篷 / 水壶 / 手电筒 / 急救包 / 食物",sz=11,c=DARK)
ap(tf_ex,"❌ 没用: 玩具 / 电视 / 大蛋糕 / 游戏机 😂",sz=11,c=DARK)
# Right: task
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.6),Inches(0.95),Inches(3.1),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=PINE;sh2.line.fill.background()
tb(s,6.7,0.98,2.9,0.35,"👉 任务  Task",sz=14,b=True,c=WHITE)
tf2=tb(s,6.7,1.45,3.0,2.3,"1️⃣ 看 9 个物品",sz=12,c=DARK)
ap(tf2,"2️⃣ 选出有用的 (✔️) 和没用的 (❌)",sz=12,c=DARK)
ap(tf2,"3️⃣ 分类:",sz=12,c=DARK)
ap(tf2,"   😴 睡觉: ___",sz=11,c=DARK)
ap(tf2,"   🍱 吃东西: ___",sz=11,c=DARK)
ap(tf2,"   🛡️ 安全: ___",sz=11,c=DARK)
# Bottom: sentence frames
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.9),Inches(9.4),Inches(1.25))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=OK;sh3.line.width=Pt(2)
tb(s,0.5,4.0,9,0.35,"🗣️ 展示句型:",sz=14,b=True,c=OK)
tb(s,0.5,4.4,9,0.35,"· 这个有用 / 没用。",sz=14,c=DARK)
tb(s,0.5,4.7,9,0.35,"· 我选 ___, 因为 ___。",sz=14,c=DARK)
tb(s,5.0,4.4,4.7,0.35,"· 睡觉需要 ___。 / 吃东西需要 ___。 / 安全需要 ___。",sz=12,c=DARK)
n+=1;pn(s,n)
notes(s,"低 level 项目: 用图卡和打勾, 不需要手工。重点是分类 + 用句型说原因。可以让学生先剪图卡再分类。")

# 48. CLOSING — mission complete
s=ns();bg(s,PINE)
tb(s,1,0.6,8,0.7,"🏆 任务完成!",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.35,8,0.5,"Mission Complete!",sz=20,c=WARM,a=PP_ALIGN.CENTER)
# Badge
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.0),Inches(3.0),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=SUNYEL;sh.line.color.rgb=WHITE;sh.line.width=Pt(4)
tb(s,3.5,2.4,3.0,0.6,"⭐",sz=80,a=PP_ALIGN.CENTER)
tb(s,3.5,3.6,3.0,0.5,"小探险家",sz=22,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,3.5,4.1,3.0,0.4,"Little Explorer",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"明天: Day 3 — 找路 + 信号 / Day 3: Navigation + Signals",sz=14,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"颁奖时间 (1 分钟): 给每个学生贴 1 张「小探险家」贴纸 / 印章。\n说: 「你今天完成了 4 个任务 — 你是真正的小探险家!」\n预告 Day 3 内容。")

# === Save ===
out="/Users/Huan/0 projects/summercourse/Chinese/野外生存与探险wilderness_pbl/day2_camp.pptx"
prs.save(out)
print(f"Saved {out}  ({n} slides)")
