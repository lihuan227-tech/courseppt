#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
野外生存与探险 Day 3 — 方向与地图 (Direction & Map · Compass Challenge)
Story-driven K-5 immersive lesson.
Session 1: Bear story → body directions → sun → natural clues → compass
Session 2: Review + 我会认 (东南西北指南针) + 我会写 (东南西北)
Session 3: Treasure Hunt + Design a City + DIY Compass
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Compass-themed palette (navy primary, red north accent) ---
PINE   = RGBColor(0x1E,0x4D,0x2B)
SUN    = RGBColor(0xE0,0x7A,0x2C)
CREAM  = RGBColor(0xFD,0xF6,0xE3)
BROWN  = RGBColor(0x6B,0x44,0x23)
SKY    = RGBColor(0x4A,0x90,0xD9)
SUNYEL = RGBColor(0xF5,0xC2,0x42)
ALERT  = RGBColor(0xD0,0x4A,0x3C)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
WARM   = RGBColor(0xFF,0xF3,0xE0)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
OK     = RGBColor(0x38,0x8E,0x3C)
NAVY    = RGBColor(0x1A,0x23,0x7E)
COMPRED = RGBColor(0xC6,0x28,0x28)
GOLD    = RGBColor(0xF9,0xA8,0x25)
N_CL = COMPRED; S_CL = NAVY; E_CL = SUN; W_CL = OK

BASE = "/Users/Huan/0 projects/summercourse/Chinese/野外生存与探险wilderness_pbl"
STORY_IMG = os.path.join(BASE, "pics", "day3_story.png")

# === Helpers ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def img(s,l,t,w,h,path,fallback="📷"):
    if os.path.exists(path):
        s.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(h))
    else:
        ib(s,l,t,w,h,fallback)
def hb(s,txt,c=NAVY,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,txt):
    s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.8,9,0.8,sub,sz=20,c=WHITE,a=PP_ALIGN.CENTER);return s
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)

def mission_card(s,l,t,w,h,num,task_cn,task_en,emoji,color):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l+0.1),Inches(t+0.08),Inches(0.55),Inches(0.55))
    badge.fill.solid();badge.fill.fore_color.rgb=color;badge.line.fill.background()
    tb(s,l+0.1,t+0.18,0.55,0.4,str(num),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+h-0.95,w-0.1,0.5,emoji,sz=30,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+h-0.45,w-0.1,0.35,task_cn,sz=14,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,l+0.7,t+0.18,w-0.8,0.4,task_en,sz=10,c=GRAY)

def teacher_student_bar(s,t,teacher_q,student_action,color=NAVY):
    """Bottom prompt bar — teacher question on left, student action on right."""
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=color;sf.line.width=Pt(2)
    tb(s,0.45,t+0.04,4.5,0.25,"👩‍🏫 老师问 Teacher asks:",sz=10,b=True,c=color)
    tb(s,0.45,t+0.27,4.5,0.28,teacher_q,sz=12,b=True,c=DARK)
    tb(s,5.0,t+0.04,4.6,0.25,"🧒 学生 Student does:",sz=10,b=True,c=SUN)
    tb(s,5.0,t+0.27,4.6,0.28,student_action,sz=12,b=True,c=DARK)

def compass_rose(s,cx,cy,radius,size_label=14):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-radius),Inches(cy-radius),Inches(radius*2),Inches(radius*2))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=NAVY;sh.line.width=Pt(2.5)
    sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-radius*0.3),Inches(cy-radius*0.3),Inches(radius*0.6),Inches(radius*0.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=GOLD;sh2.line.color.rgb=NAVY;sh2.line.width=Pt(1.5)
    ln1=s.shapes.add_connector(1,Inches(cx-radius*0.95),Inches(cy),Inches(cx+radius*0.95),Inches(cy))
    ln1.line.color.rgb=LGRAY;ln1.line.width=Pt(0.8)
    ln2=s.shapes.add_connector(1,Inches(cx),Inches(cy-radius*0.95),Inches(cx),Inches(cy+radius*0.95))
    ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.8)
    tb(s,cx-0.3,cy-radius-0.45,0.6,0.4,"北 N",sz=size_label,b=True,c=COMPRED,a=PP_ALIGN.CENTER)
    tb(s,cx-0.3,cy+radius+0.05,0.6,0.4,"南 S",sz=size_label,b=True,c=NAVY,a=PP_ALIGN.CENTER)
    tb(s,cx+radius+0.05,cy-0.2,0.7,0.4,"东 E",sz=size_label,b=True,c=SUN,a=PP_ALIGN.CENTER)
    tb(s,cx-radius-0.75,cy-0.2,0.7,0.4,"西 W",sz=size_label,b=True,c=OK,a=PP_ALIGN.CENTER)

def ab_slide(title_cn,title_en,question_cn,question_en,a_emoji,a_label,a_caption,b_emoji,b_label,b_caption,answer,reason):
    s=ns();bg(s,CREAM);hb(s,f"🧭 {title_cn}  {title_en}",NAVY)
    tb(s,0.4,0.85,9.2,0.4,question_cn,sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.20,9.2,0.3,question_en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    a_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.6),Inches(4.5),Inches(2.4))
    a_box.fill.solid();a_box.fill.fore_color.rgb=WHITE;a_box.line.color.rgb=NAVY;a_box.line.width=Pt(2.5)
    pill(s,0.5,1.7,0.7,0.4,"A",NAVY,sz=16)
    tb(s,1.2,1.65,3.6,0.5,a_label,sz=18,b=True,c=NAVY)
    tb(s,0.6,2.2,4.2,0.5,a_emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,0.6,3.1,4.2,0.6,a_caption,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    b_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.6),Inches(4.5),Inches(2.4))
    b_box.fill.solid();b_box.fill.fore_color.rgb=WHITE;b_box.line.color.rgb=ALERT;b_box.line.width=Pt(2.5)
    pill(s,5.2,1.7,0.7,0.4,"B",ALERT,sz=16)
    tb(s,5.9,1.65,3.6,0.5,b_label,sz=18,b=True,c=ALERT)
    tb(s,5.3,2.2,4.2,0.5,b_emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,5.3,3.1,4.2,0.6,b_caption,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    teacher_student_bar(s,4.20,
        "你选 A 还是 B? 走到那边!",
        "我选 ___, 因为 ___ 。 (走到 A 边 / B 边)")
    tb(s,0.4,4.85,9.2,0.3,"💬 I choose A/B because…",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    notes(s,f"老师备课:\n• 答案: {answer}\n• 原因: {reason}")
    return s

def word_card_read(w,py,en,img_emoji,sent,color=NAVY):
    """Image-driven recognition — emoji/picture LEFT, character + pinyin RIGHT, example sentence on bottom."""
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    # Left: image cue
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(2.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,0.5,1.3,4.2,2.0,img_emoji,sz=130,a=PP_ALIGN.CENTER)
    # Right: big character — shrink for multi-char words
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(2.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=color;sh2.line.width=Pt(2.5)
    char_sz = 110 if len(w)==1 else (66 if len(w)==2 else 54)
    tb(s,5.1,1.15,4.4,1.5,w,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,5.1,2.70,4.4,0.4,f"{py}  ·  {en}",sz=18,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,5.1,3.15,4.4,0.4,"👉 跟我读 3 次！",sz=14,c=color,a=PP_ALIGN.CENTER)
    # Bottom: example sentence
    sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.3))
    sh3.fill.solid();sh3.fill.fore_color.rgb=WHITE;sh3.line.color.rgb=color;sh3.line.width=Pt(2)
    tb(s,0.6,3.95,2.0,0.4,"💬 例句 Example",sz=14,b=True,c=color)
    tb(s,0.6,4.35,8.8,0.6,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,strokes_count,strokes_hint,color=NAVY):
    """Watch → air-write → write."""
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    # Left: big char + pinyin
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,0.5,1.1,4.2,2.4,w,sz=160,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.55,4.2,0.5,f"{py}  ·  {en}",sz=22,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.10,4.2,0.4,f"{strokes_count} 笔 / {strokes_count} strokes",sz=16,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.50,4.2,0.4,strokes_hint,sz=11,c=DARK,a=PP_ALIGN.CENTER)
    # Right top: 3-step practice
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,5.15,1.1,4.4,0.4,"3 步练习  3 Steps",sz=16,b=True,c=color)
    tb(s,5.15,1.5,4.4,0.4,"1️⃣ 看老师写  Watch teacher",sz=13,c=DARK)
    tb(s,5.15,1.85,4.4,0.4,"2️⃣ 用手指空中写  Air-write",sz=13,c=DARK)
    tb(s,5.15,2.20,4.4,0.4,"3️⃣ 在田字格写 3 次  Write 3 times",sz=13,c=DARK)
    # Right bottom: 4 田字格 squares for practice
    for i in range(4):
        x=5.0+i*1.15
        sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.85),Inches(1.05),Inches(1.05))
        sq.fill.solid();sq.fill.fore_color.rgb=WHITE;sq.line.color.rgb=color;sq.line.width=Pt(1.5)
        ln1=s.shapes.add_connector(1,Inches(x),Inches(2.85+0.525),Inches(x+1.05),Inches(2.85+0.525))
        ln1.line.color.rgb=LGRAY;ln1.line.width=Pt(0.5);ln1.line.dash_style=2
        ln2=s.shapes.add_connector(1,Inches(x+0.525),Inches(2.85),Inches(x+0.525),Inches(2.85+1.05))
        ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.5);ln2.line.dash_style=2
    tb(s,5.0,4.0,4.6,0.3,"在田字格里写 3 次 ↓  Write here 3 times",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    teacher_student_bar(s,4.45,
        f"和我一起写「{w}」",
        "看 → 空中写 → 田字格写 3 次")
    return s

# ========================================================================
#                              SLIDES
# ========================================================================
n=0

# 1. COVER
s=ns();bg(s,NAVY)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.4,8,0.5,"DAY 3",sz=18,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"🧭 方向与地图",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"Direction & Map  ·  指南针挑战 Compass Challenge",sz=20,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"🧭 探险家任务  Explorer Mission",sz=24,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"听故事 · 找方向 · 做指南针 · 画地图 · 寻宝",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"Story · Find · Make · Map · Hunt",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"野外生存与探险 · Wilderness Survival",sz=14,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「探险家们! 今天有一只小熊迷路了 — 我们要帮它回家!」\n• 不讲方向是什么 — 让学生先听故事再发现答案。")

# 2. SESSION 1 LEARNING GOALS
s=ns();bg(s,CREAM);hb(s,"🎯 Session 1 学习目标  Today's Goals",NAVY)
tb(s,0.4,0.85,9.2,0.4,"上完 Session 1, 你可以…",sz=18,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"After Session 1, you can…",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
goals=[
    ("1️⃣","在地图上认识 4 个方向","Identify 4 directions on a map: 东南西北",N_CL),
    ("2️⃣","用太阳找方向","Use the sun (morning = east)",SUN),
    ("3️⃣","看大自然的线索","Read nature's clues (sun · trees · moss · ants)",PINE),
    ("4️⃣","认识指南针","Know a compass (red needle = North)",COMPRED),
    ("5️⃣","听指令走方向","Follow direction instructions (e.g., 向北 3 步)",NAVY),
]
for i,(num,cn,en,cl) in enumerate(goals):
    y=1.70+i*0.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.58))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.55,y+0.10,0.7,0.4,num,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.05,3.5,0.3,cn,sz=17,b=True,c=cl)
    tb(s,5.00,y+0.10,4.5,0.3,en,sz=12,c=DARK)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「上 Session 1 我们要学会 5 件事 — 一起读!」\n• 跟读 5 个目标, 完成 1 个打 1 个 ✓。\n• 关键词: 4 方向 / 太阳 / 大自然 / 指南针 / 听指令。")

# 3. SESSION 1 DIVIDER
s=div("Session 1  上午 50 min","🐻 故事 · 方向 · 太阳 · 指南针  Story-driven exploration",NAVY,"📖");n+=1;pn(s,n)

# 4. STORY — Hook with full image
s=ns();bg(s,CREAM);hb(s,"🐻 故事时间  Story Time",NAVY)
img(s,0.3,0.95,9.4,3.6,STORY_IMG,"📷 Bear lost story illustration")
tb(s,0.4,4.65,9.2,0.4,"小熊去露营。它走到河边, 迷路了。",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,5.05,9.2,0.3,"A little bear goes camping. It walks to the river and gets lost.",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"故事打开 (2 分钟):\n• 慢慢念故事, 看图。\n• 不要急着给答案 — 让学生看图想象。\n• 「小熊看到了什么? 它的帐篷在哪里?」")

# 5. STORY — How can the bear find the way? (QUESTION ONLY — students think first)
s=ns();bg(s,CREAM);hb(s,"🤔 怎么办?  How can the bear find the way?",SUN)
tb(s,0.4,0.95,9.2,0.5,"想一想 — 小熊可以用什么找路?",sz=24,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.45,9.2,0.3,"Think — what can the bear use to find the way?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# 4 question-mark cards (no answer revealed)
for i in range(4):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(2.20),Inches(2.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=NAVY;sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.35,2.1,1.0,"❓",sz=80,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,2.1,0.4,f"想法 {i+1}",sz=16,b=True,c=NAVY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.95,2.1,0.3,f"Idea {i+1}",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,
    "和你旁边的同学说一说!",
    "Think-Pair-Share: 我觉得小熊可以用 ___ 。")
n+=1;pn(s,n)
notes(s,"think-pair-share (3 分钟):\n• 不要给答案!\n• 让学生 think → pair → share, 收集所有答案在白板上。\n• 等学生说完, 再翻到下一页揭晓。")

# 6. STORY — Answer reveal (after students share)
s=ns();bg(s,CREAM);hb(s,"💡 我们的好办法!  Our Ideas!",SUN)
tb(s,0.4,0.85,9.2,0.4,"你说得对! 这些都可以帮小熊!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"You got it! These can all help the bear.",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
hints=[("☀️","看太阳","Look at the sun",SUN),
       ("🌳","看树","Look at trees",PINE),
       ("🧭","用指南针","Use a compass",COMPRED),
       ("🗺️","看地图","Look at a map",NAVY)]
for i,(em,cn,en,cl) in enumerate(hints):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.20),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.95,2.1,1.0,em,sz=64,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,2.1,0.5,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,2.1,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,
    "今天我们都试试这些方法!",
    "一起说: 太阳! 树! 指南针! 地图!")
n+=1;pn(s,n)
notes(s,"答案揭晓 (1 分钟):\n• 「你刚才说的, 是不是这 4 个?」\n• 表扬学生想到的新答案 (例: 看星星, 听风, 等等)\n• 「今天我们重点学这 4 个 — 准备好了吗?」")

# 7. ON THE MAP — body to map (image LEFT, ordered rows RIGHT)
s=ns();bg(s,CREAM);hb(s,"🗺️ 在地图上辨别方向  On the Map",NAVY)
tb(s,0.4,0.85,9.2,0.4,"上北 · 下南 · 左西 · 右东",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Up=North · Down=South · Left=West · Right=East",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# LEFT: image placeholder for body / compass illustration
ib(s,0.4,1.65,3.8,2.95,"📷 小朋友面向北的图  Person facing north (body→NSEW)")
# RIGHT: 4 mapping rows in 上北 → 下南 → 左西 → 右东 order
mappings=[("👆 上","头顶","Head","北 N",N_CL,1.65),
          ("👇 下","脚下","Feet","南 S",S_CL,2.40),
          ("👈 左","左手","Left","西 W",W_CL,3.15),
          ("👉 右","右手","Right","东 E",E_CL,3.90)]
for em_l,bp_cn,bp_en,nsew,cl,y in mappings:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.4),Inches(y),Inches(5.2),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,4.55,y+0.13,1.4,0.4,em_l,sz=20,b=True,c=cl)
    tb(s,5.95,y+0.06,1.6,0.4,bp_cn,sz=14,b=True,c=DARK)
    tb(s,5.95,y+0.36,1.6,0.3,bp_en,sz=9,c=GRAY)
    tb(s,7.55,y+0.13,0.7,0.4,"→",sz=22,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,8.25,y+0.13,1.3,0.4,nsew,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.65,
    "头顶是哪边? 脚下呢?",
    "全班站起来, 用手指: 上北! 下南! 左西! 右东!")
n+=1;pn(s,n)
notes(s,"互动 (3 分钟):\n• 全班站起来, 跟着老师指方向。\n• 老师面向北, 学生跟着指: 上 (头顶) 北! 下 (脚下) 南! 右手东! 左手西!\n• 重复 3 次, 越来越快。")

# 8. MAP RULES — 上北下南左西右东 (with camp map illustration)
s=ns();bg(s,CREAM);hb(s,"🗺️ 地图的规则  Map Rules",NAVY)
# Big illustration centered (image's own title says 上北下南，左西右东是地图的规则)
img(s,1.5,0.95,7.0,3.65,os.path.join(BASE,"pics","day 3 slide 8.png"),
    "📷 地图规则 / Map rules illustration (上北下南左西右东)")
teacher_student_bar(s,4.70,
    "帐篷在哪边? 火在哪边? 大树在哪边?",
    "用手指: 帐篷在北边, 火在西边, 大树在东边。")
n+=1;pn(s,n)
notes(s,"K-5 关键 (3 分钟):\n• 看图 — 地图最上面是「北」\n• 一起念: 上北 — 下南 — 左西 — 右东\n• 让学生用手指图上的物品: 帐篷在哪边? 火堆? 大树?\n• 「记住这 8 个字, 你就能看懂地图啦!」\n• 这是后面 Project 2 (Design My City) 的基础")

# 8b. COMMUNITY MAP — Basic (level 1): find direction
s=ns();bg(s,CREAM);hb(s,"🗺️ 我们的小区地图  Our Community Map",NAVY)
# Top: difficulty pill (wider so text fits on one line)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.80),Inches(4.4),Inches(0.40))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.fill.background()
tb(s,0.4,0.85,4.4,0.32,"🟡 基础 — 找方向 Find Direction",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# LEFT: community map image
img(s,0.4,1.30,5.4,3.40,os.path.join(BASE,"pics","D3 community map.png"),"📷 Community map (邮局/学校/泳池/公园/消防局/商场/图书馆/池塘)")
# RIGHT: 8 questions (2 columns of 4)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.0),Inches(1.30),Inches(3.7),Inches(3.40))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=NAVY;sh2.line.width=Pt(2)
tb(s,6.10,1.40,3.5,0.30,"❓ ___ 在哪个方向?",sz=13,b=True,c=NAVY)
tb(s,6.10,1.72,3.5,0.25,"Which direction is ___?",sz=9,c=GRAY)
qs=["📮 邮局","🏫 学校","🏊 游泳池","🌳 公园","🚒 消防局","🛍️ 商场","📚 图书馆","🪷 池塘"]
for i,q in enumerate(qs):
    col=i%2;row=i//2
    x=6.10+col*1.85;y=2.10+row*0.58
    tb(s,x,y,1.7,0.4,q,sz=13,b=True,c=DARK)
teacher_student_bar(s,4.80,
    "用手指地图! ___ 在哪个方向?",
    "「___ 在 ___ 边。」(北/南/西/东/东北/...)")
n+=1;pn(s,n)
notes(s,"基础 (5-6 分钟):\n• 让学生用手指地图, 一个一个回答\n• 答案 (上北下南左西右东规则下):\n  - 邮局: 东边 (中排东)\n  - 学校: 北边 (上排中)\n  - 游泳池: 东北 (上排东)\n  - 公园: 西边 (中排西)\n  - 消防局: 中间\n  - 商场: 南边 (下排中)\n  - 图书馆: 东南 (下排东)\n  - 池塘: 西南 (下排西)\n• 句型: 「邮局在东边。」")

# 8c. COMMUNITY MAP — Compare (level 2) + Challenge (level 3)
s=ns();bg(s,CREAM);hb(s,"🗺️ 小区地图 — 进阶  Community Map · Advanced",NAVY)
# LEFT: same map (smaller)
img(s,0.3,0.85,4.4,3.0,os.path.join(BASE,"pics","D3 community map.png"),"📷 Community map")
# RIGHT TOP: comparison questions
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.85),Inches(4.85),Inches(0.42))
sh.fill.solid();sh.fill.fore_color.rgb=COMPRED;sh.line.fill.background()
tb(s,4.95,0.90,4.65,0.35,"🟠 比较位置  Compare Positions (稍难)",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.35),Inches(4.85),Inches(2.50))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=COMPRED;sh2.line.width=Pt(2)
compare_qs=[
    "📮 邮局在学校的哪一边?",
    "🛍️ 商场在公园的哪一边?",
    "📚 图书馆在商场的哪一边?",
    "🚒 消防局在公园的哪一边?",
]
for i,q in enumerate(compare_qs):
    tb(s,5.0,1.45+i*0.55,4.6,0.4,f"{i+1}. {q}",sz=13,b=True,c=DARK)
# BOTTOM (across full width): challenge questions
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.0),Inches(9.4),Inches(0.42))
sh3.fill.solid();sh3.fill.fore_color.rgb=BROWN;sh3.line.fill.background()
tb(s,0.4,4.05,9.2,0.35,"🟤 挑战题 G2+  Challenge (G2+)",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
sh4=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.50),Inches(9.4),Inches(0.55))
sh4.fill.solid();sh4.fill.fore_color.rgb=WARM;sh4.line.color.rgb=BROWN;sh4.line.width=Pt(2)
tb(s,0.45,4.57,4.65,0.4,"1. 从学校出发, 往东走 — 你会看到什么?",sz=12,b=True,c=DARK)
tb(s,5.10,4.57,4.55,0.4,"2. 从公园出发, 往北走 — 会到哪里?",sz=12,b=True,c=DARK)
teacher_student_bar(s,5.10,
    "想一想 — 用手指走路!",
    "「___ 在 ___ 的 ___ 边。」/ 「往 ___ 走, 我看见 ___ 。」")
n+=1;pn(s,n)
notes(s,"进阶 (4-5 分钟):\n• 比较位置答案:\n  1. 邮局在学校的东南 / 右下边 (邮局: 中排东; 学校: 上排中)\n  2. 商场在公园的东南边 (商场: 下排中; 公园: 中排西)\n  3. 图书馆在商场的东边 (都在下排)\n  4. 消防局在公园的东边 (都在中排)\n• 挑战题答案:\n  1. 从学校 (上排中) 往东 → 游泳池 (上排东)\n  2. 从公园 (中排西) 往北 → 房子 / HOUSE (上排西)\n• 让学生用手在地图上「走」, 一格一格指")

# 9. SUN → EAST (multi-choice question + tape "东" label)
s=ns();bg(s,CREAM);hb(s,"☀️ 太阳从哪里升起?  Where does the sun rise?",SUN)
# Top: multi-choice question (4 options as colored cards)
tb(s,0.4,0.85,9.2,0.4,"❓ 你知道太阳从哪里升起吗?",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Do you know where the sun rises?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
opts=[("A","🔝 北","North",N_CL),
      ("B","➡️ 东","East",E_CL),
      ("C","⬇️ 南","South",S_CL),
      ("D","⬅️ 西","West",W_CL)]
for i,(letter,cn,en,cl) in enumerate(opts):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(2.20),Inches(1.40))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    pill(s,x+0.1,1.80,0.55,0.40,letter,cl,sz=14)
    tb(s,x+0.1,2.15,2.0,0.5,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.65,2.0,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Action sequence (3 steps below)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.30),Inches(9.2),Inches(1.20))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2)
tb(s,0.6,3.40,2.9,0.4,"1️⃣ 学生选答案",sz=13,b=True,c=SUN)
tb(s,0.6,3.75,2.9,0.4,"举手: A / B / C / D",sz=14,b=True,c=DARK)
tb(s,3.6,3.40,3.0,0.4,"2️⃣ 用手指太阳",sz=13,b=True,c=SUN)
tb(s,3.6,3.75,3.0,0.4,"指窗外 — 那边!",sz=14,b=True,c=DARK)
tb(s,6.7,3.40,2.9,0.4,"3️⃣ 老师贴标签",sz=13,b=True,c=SUN)
tb(s,6.7,3.75,2.9,0.4,"在墙上贴「东」",sz=14,b=True,c=E_CL)
# Big east label
badge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.10),Inches(9.2),Inches(0.4))
badge.fill.solid();badge.fill.fore_color.rgb=E_CL;badge.line.fill.background()
tb(s,0.4,4.13,9.2,0.35,"✅ 答案: B 东 East — 那一面墙 = 东",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.65,
    "举手! A B C D — 你选哪个?",
    "选答案 → 指太阳 → 看老师贴「东」标签!")
n+=1;pn(s,n)
notes(s,"探究式开场 (4-5 分钟):\n• Step 1: 提问 + 4 个选项 — 学生举手投票\n• 不要立刻说答案! 让学生先猜\n• Step 2: 让所有学生用手指真的指向窗外 (太阳那边)\n• Step 3: 老师拿一张写有「东」的 A4 纸, 走过去贴在那一面墙 — 这个动作是关键! \n• 这就锚定了 east, 后面儿歌 + 活动都是从这里开始\n• 接下来 (next slide): 学方向儿歌")

# 10. CHANT + BODY MOVEMENT
s=ns();bg(s,CREAM);hb(s,"🎵 方向儿歌 + 身体动作  Chant + Body",SUN)
# Action steps strip on top
acts=[("🧍","站起来","Stand up",NAVY),
      ("👀","面向东墙","Face EAST wall",E_CL),
      ("🗣️","边读边指","Chant + point",COMPRED)]
for i,(em,cn,en,cl) in enumerate(acts):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(3.0),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,x+0.1,1.05,0.6,0.5,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.75,1.00,2.2,0.4,cn,sz=15,b=True,c=cl)
    tb(s,x+0.75,1.30,2.2,0.3,en,sz=10,c=GRAY)
# Chant box (full-width)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.80),Inches(9.2),Inches(2.55))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=SUN;sh3.line.width=Pt(2.5)
# 3 chant lines with action arrows
tb(s,0.6,1.95,9.0,0.5,"早晨起来, 面向太阳。",sz=24,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Line 2: 前面是东 ↔ 后面是西 (split layout: chant left + actions right)
tb(s,0.6,2.55,4.4,0.5,"前面是东,",sz=22,b=True,c=E_CL)
tb(s,5.0,2.55,4.4,0.5,"后面是西。",sz=22,b=True,c=W_CL)
tb(s,0.6,3.00,4.4,0.3,"👉 手往前指 → 东",sz=12,b=True,c=GRAY)
tb(s,5.0,3.00,4.4,0.3,"👉 手往后指 → 西",sz=12,b=True,c=GRAY)
# Line 3: 左面是北 ↔ 右面是南
tb(s,0.6,3.40,4.4,0.5,"左面是北,",sz=22,b=True,c=N_CL)
tb(s,5.0,3.40,4.4,0.5,"右面是南。",sz=22,b=True,c=S_CL)
tb(s,0.6,3.85,4.4,0.3,"👈 左手 → 北",sz=12,b=True,c=GRAY)
tb(s,5.0,3.85,4.4,0.3,"👉 右手 → 南",sz=12,b=True,c=GRAY)
teacher_student_bar(s,4.50,
    "面向东! 准备! 1, 2, 3 — 读!",
    "全班站起来, 面向东墙, 边读边用手指!")
n+=1;pn(s,n)
notes(s,"互动儿歌 (5 分钟):\n• Step 1: Stand → Face East → Chant\n• 老师领读 1 次, 全班跟读 3 次\n• 重点动作:\n  - 「前面是东」 — 双手往前指\n  - 「后面是西」 — 双手往后指 (转头看后面)\n  - 「左面是北」 — 左手往左伸\n  - 「右面是南」 — 右手往右伸\n• 一定要让学生先面向「东」墙 (slide 9 贴的标签那边)\n• 越读越快, 最后让学生背下来")

# 11. FIND THE PATTERN — clockwise turn game (now BEFORE follow-commands)
s=ns();bg(s,CREAM);hb(s,"🔄 你发现了什么?  Find the Pattern!",E_CL)
# Action prompt
tb(s,0.4,0.85,9.2,0.4,"全班面向: 东 → 南 → 西 → 北 → 东",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Face: East → South → West → North → East",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Compass cycle: 4 direction circles in a clockwise arrow flow
dirs=[("东","E",E_CL,1.0),("南","S",S_CL,3.05),("西","W",W_CL,5.10),("北","N",N_CL,7.15)]
for cn,en,cl,x in dirs:
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(1.75),Inches(1.7),Inches(1.7))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.color.rgb=WHITE;sh.line.width=Pt(3)
    tb(s,x,1.80,1.7,0.7,cn,sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x,2.65,1.7,0.5,en,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# Arrows between circles
for ax in [2.75,4.80,6.85]:
    tb(s,ax,2.25,0.30,0.6,"→",sz=32,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Loop arrow back to 东
tb(s,8.85,2.25,0.6,0.6,"↻",sz=36,b=True,c=GOLD,a=PP_ALIGN.CENTER)
# Pattern findings (2 cards, side by side)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.60),Inches(4.5),Inches(0.85))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=GOLD;sh.line.width=Pt(2)
tb(s,0.55,3.65,4.3,0.30,"💡 一直向右转!  Like a clock ⏰",sz=13,b=True,c=GOLD)
tb(s,0.55,3.98,4.3,0.42,"像时钟一样!",sz=20,b=True,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(3.60),Inches(4.5),Inches(0.85))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=COMPRED;sh2.line.width=Pt(2)
tb(s,5.25,3.65,4.3,0.30,"🪞 对面!  Opposites",sz=13,b=True,c=COMPRED)
tb(s,5.25,3.98,4.3,0.42,"东 ↔ 西  ·  南 ↔ 北",sz=20,b=True,c=DARK)
teacher_student_bar(s,4.75,
    "知道东, 你能找到南、西、北吗?",
    "右转 → 南! 再右转 → 西! 再右转 → 北!")
n+=1;pn(s,n)
notes(s,"动起来 + 发现规律 (5 分钟):\n• Step 1: 全班面向东\n• Step 2: 老师喊 「向右转!」 → 学生转 90° → 「这是哪边?」 (南)\n• Step 3: 再喊「向右转!」 → 西 → 再转 → 北 → 再转 → 回到东\n• 关键问题: 「你发现了什么?」 (Don't tell — let them notice)\n• 期望答案: 「一直向右转」「像时钟一样」「东西相对, 南北相对」\n• Challenge: 「我说东 — 你能找到南/西/北吗?」 — 全班一起转")

# 11b. INTERACTIVE — student volunteer follows multi-step direction commands (now AFTER pattern)
s=ns();bg(s,CREAM);hb(s,"🎮 跟着指令走!  Follow the Commands!",E_CL)
tb(s,0.4,0.85,9.2,0.4,"老师找 1 位同学上来 — 听指令走!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Teacher invites a volunteer — follow the commands!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Sample command cards (3 examples) — left side
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.70),Inches(5.2),Inches(2.95))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=E_CL;sh.line.width=Pt(2.5)
tb(s,0.55,1.80,5.0,0.4,"📢 指令例子  Example Commands",sz=14,b=True,c=E_CL)
cmds=[
    ("1️⃣","往北走 3 步","→ 3 steps NORTH",N_CL),
    ("2️⃣","往南走 5 步","→ 5 steps SOUTH",S_CL),
    ("3️⃣","往东走 2 步","→ 2 steps EAST",E_CL),
    ("4️⃣","往西走 4 步","→ 4 steps WEST",W_CL),
]
for i,(num,cn,en,cl) in enumerate(cmds):
    y=2.25+i*0.55
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(y),Inches(4.9),Inches(0.50))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=cl;sh2.line.width=Pt(1.5)
    tb(s,0.65,y+0.10,0.55,0.35,num,sz=18,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.06,2.4,0.4,cn,sz=15,b=True,c=cl)
    tb(s,3.75,y+0.10,1.7,0.3,en,sz=11,c=GRAY)
# RIGHT: how-to-play steps
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.85),Inches(1.70),Inches(3.85),Inches(2.95))
sh3.fill.solid();sh3.fill.fore_color.rgb=WHITE;sh3.line.color.rgb=NAVY;sh3.line.width=Pt(2.5)
tb(s,6.0,1.80,3.7,0.4,"🎯 怎么玩  How to Play",sz=14,b=True,c=NAVY)
tf=tb(s,6.0,2.25,3.7,0.4,"1. 老师找 1 位志愿者",sz=12,b=True,c=DARK)
ap(tf,"2. 学生站在教室中间",sz=12,b=True,c=DARK)
ap(tf,"3. 面向「东」 (东墙)",sz=12,b=True,c=E_CL)
ap(tf,"4. 听老师指令 + 走!",sz=12,b=True,c=DARK)
ap(tf,"5. 全班喊: 「对!」/「不对!」",sz=12,b=True,c=COMPRED)
ap(tf,"6. 换 3-4 位 volunteer",sz=12,b=True,c=DARK)
teacher_student_bar(s,4.75,
    "志愿者: 往北走 3 步! 往南走 5 步! ...",
    "面向东 → 听 → 走! 全班帮忙数 1, 2, 3!")
n+=1;pn(s,n)
notes(s,"互动游戏 (6-8 分钟):\n• 找 1 位 volunteer 上来站在教室中间\n• 让 volunteer 先面向「东」墙 (slide 9 贴的标签)\n• 老师喊指令, volunteer 走 — 全班一起数步数\n• 例:\n  - 「往北走 3 步」 → volunteer 转身向北走 3 步\n  - 「往南走 5 步」 → 转身向南走 5 步\n  - 「往东走 2 步」 → 转身向东走 2 步\n• 全班监督: 走错了喊「不对!」, 走对了喊「对!」\n• 进阶: 「往东北走 2 步」(8 个方向)\n• 换 3-4 位 volunteer 各做一次 (每次 3-5 个指令)\n• 这是 Project 1 (Treasure Hunt) 的预热")

# 12-15. NATURAL CLUES — one slide per clue (image LEFT, fact + caution RIGHT)
def nature_clue_slide(em,cn,en,fact_cn,fact_en,caution_cn,caution_en,img_hint,color):
    s=ns();bg(s,CREAM);hb(s,f"🌿 大自然的线索 — {em} {cn}  Nature's Clue: {en}",color)
    # LEFT: image placeholder
    ib(s,0.4,1.0,4.4,3.6,f"📷 {img_hint}")
    # RIGHT: 2 stacked cards (fact + caution)
    # Fact card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,5.15,1.10,4.4,0.4,f"{em}  小秘密",sz=15,b=True,c=color)
    tb(s,5.15,1.55,4.4,0.6,fact_cn,sz=22,b=True,c=DARK)
    tb(s,5.15,2.20,4.4,0.4,fact_en,sz=12,c=GRAY)
    tb(s,5.15,2.55,4.4,0.4,"💬 我看 ___, 那边可能是 ___ 。",sz=12,b=True,c=color)
    # Caution card
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(3.10),Inches(4.6),Inches(1.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=ALERT;sh2.line.width=Pt(2)
    tb(s,5.15,3.20,4.4,0.35,"⚠️ 注意 Caution",sz=13,b=True,c=ALERT)
    tb(s,5.15,3.55,4.4,0.4,caution_cn,sz=14,b=True,c=DARK)
    tb(s,5.15,3.95,4.4,0.4,caution_en,sz=10,c=GRAY)
    tb(s,5.15,4.25,4.4,0.3,"➜ 最准的是指南针! Most reliable: compass",sz=11,b=True,c=COMPRED)
    teacher_student_bar(s,4.75,
        f"{cn} 告诉我们什么?",
        f"用手指 — 我看 {cn}, 那边是 ___ !")
    return s

# 12. SUN
s=nature_clue_slide("☀️","太阳","Sun",
    "早上→东 · 下午→西","Morning→East · Evening→West",
    "中午太阳在头顶, 不容易判断","At noon, sun is overhead — harder to tell",
    "早上太阳从东边升起",SUN)
n+=1;pn(s,n)
notes(s,"☀️ 太阳 (2 分钟):\n• 这是最简单, 也是最常用的方法\n• 让学生回忆: 早上太阳在哪? 下午呢?\n• 配合 slide 10 儿歌一起复习\n• 局限: 阴天看不到太阳 / 中午太阳在头顶")

# 13. TREE LEAVES
s=nature_clue_slide("🌳","树叶","Tree Leaves",
    "茂盛的一面 = 朝阳 (常朝南)","Lusher side = sunnier (often south)",
    "城市里树被建筑挡住就不准","In cities, buildings block sun — unreliable",
    "树的两面叶子稠密对比",PINE)
n+=1;pn(s,n)
notes(s,"🌳 树叶 (2 分钟):\n• 阳光多的那一面叶子长得茂\n• 在北半球, 朝南 (向阳) 那一面更茂盛\n• 局限: 城市里树两边光照都受建筑影响, 不准")

# 14. MOSS
s=nature_clue_slide("🍀","青苔","Moss",
    "常长在树的「北」面","Often grows on the NORTH side of trees",
    "潮湿的地方四面都会长", "In wet places, moss grows everywhere",
    "树干北面长青苔特写",OK)
n+=1;pn(s,n)
notes(s,"🍀 青苔 (2 分钟):\n• 北面背阴 + 潮湿 → 青苔喜欢\n• 让学生想象: 阳光晒哪一面? 不晒哪一面? (背阴 = 北)\n• 局限: 在很潮湿的森林里, 四面都有青苔")

# 15. ANT HILLS
s=nature_clue_slide("🐜","蚁丘","Ant Hills",
    "开口常朝南 (避风向阳)","Opening often faces SOUTH (warm + sheltered)",
    "蚂蚁有时也会换方向","Ants sometimes change direction",
    "蚂蚁丘开口",BROWN)
n+=1;pn(s,n)
notes(s,"🐜 蚁丘 (2 分钟):\n• 蚂蚁喜欢温暖避风的家\n• 北半球阳光从南来 → 开口朝南\n• 局限: 不是 100% 规则, 蚂蚁会根据地形调整")

# 16. TEACHER VIDEO — natural clues
s=ns();bg(s,NAVY)
tb(s,1,0.30,8,0.6,"🎬 看视频  Watch a Video",sz=28,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,0.90,8,0.35,"大自然的方向线索 / Nature's Direction Clues",sz=15,b=True,c=WARM,a=PP_ALIGN.CENTER)
# Larger video placeholder (16:9 frame, near full-width)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.40),Inches(9.0),Inches(3.50))
sh.fill.solid();sh.fill.fore_color.rgb=DARK;sh.line.color.rgb=WHITE;sh.line.width=Pt(3)
tb(s,0.5,2.30,9.0,1.5,"▶️",sz=160,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,4.40,9.0,0.40,"🔗 老师在这里插入视频 / Teacher: paste video here",sz=14,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"老师准备:\n• 提前找 1 个 K-5 友好的视频 (YouTube 搜「nature direction clues for kids」「找方向自然」)\n• 推荐: National Geographic Kids / SciShow Kids 类型\n• 视频长度: 2-4 分钟最佳\n• 看视频前布置任务: 数有几个自然线索\n• 看视频后总结: 自然线索可以参考, 但不100% 准 — 指南针最准!")

# 13. COMPASS — DISCOVERY ACTIVITY (hand out compasses, observe red needle)
s=ns();bg(s,CREAM);hb(s,"🧭 一起玩指南针!  Let's Play with the Compass!",COMPRED)
tb(s,0.4,0.85,9.2,0.4,"老师给每位小朋友一个指南针 — 大家一起找一个神奇的秘密!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.3,"Each child gets a compass — let's discover something magic!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 4 step cards
steps=[("1️⃣","拿一个指南针","Hold a compass","🧭",COMPRED),
       ("2️⃣","站到教室不同地方","Stand in different spots","🚶",NAVY),
       ("3️⃣","平拿 ✋ 等它停","Hold flat — wait for it to settle","✋",E_CL),
       ("4️⃣","看! 红针指哪里?","Look! Red needle points where?","🔴",ALERT)]
for i,(num,cn,en,em,cl) in enumerate(steps):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.20),Inches(2.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.75,2.0,0.4,num,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.20,2.0,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.10,2.0,0.4,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.55,2.0,0.4,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Discovery hint
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.20),Inches(9.2),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.fill.background()
tb(s,0.4,4.27,9.2,0.4,"🤔 红针都指同一面墙吗?  Does the red needle always point to the SAME wall?",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.80,
    "你站在哪里? 红针指哪面墙?",
    "我站在 ___, 红针指 ___ 面墙。")
n+=1;pn(s,n)
notes(s,"探究活动 (5-7 分钟):\n• 老师准备每位学生 1 个小指南针 (或每组 1 个)\n• 让学生拿着指南针, 站到教室不同位置, 面向不同方向\n• 关键: 平拿 ✋ — 不要斜着\n• 等 1-2 秒, 让针停下来\n• 老师巡视: 「你的红针指哪面墙?」\n• 不要立刻给答案 — 让学生自己发现\n• 集中 share: 几位学生回答位置 + 红针指向\n• 引导发现: 不管站在哪, 红针都指同一面墙!\n• 这正好是 slide 11 学过的「东」墙的对面 = 北墙")

# 14. COMPASS — REVEAL: red needle always points NORTH
s=ns();bg(s,CREAM);hb(s,"🔴 红针永远指北!  Red Needle Always Points NORTH!",COMPRED)
# LEFT: image placeholder (teacher inserts compass photo) + anatomy callouts below
tb(s,0.5,1.05,4.2,0.4,"🧭 指南针长这样:",sz=15,b=True,c=COMPRED,a=PP_ALIGN.CENTER)
ib(s,0.4,1.50,4.4,2.50,"📷 老师插入指南针图片 / Teacher: insert compass image")
# Anatomy callouts
tb(s,0.55,4.10,4.2,0.35,"🎯 方位盘 = 圆圆的盘 (东南西北)",sz=12,b=True,c=NAVY)
tb(s,0.55,4.40,4.2,0.35,"📍 磁针 = 中间会动的针 (红 = 北!)",sz=12,b=True,c=COMPRED)
# RIGHT: key principle
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(0.95),Inches(4.6),Inches(3.7))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=COMPRED;sh2.line.width=Pt(2.5)
tb(s,5.15,1.10,4.4,0.4,"💡 神奇的秘密  The Magic:",sz=15,b=True,c=COMPRED)
tb(s,5.15,1.55,4.4,0.6,"不论你站在哪里…",sz=20,b=True,c=DARK)
tb(s,5.15,2.10,4.4,0.6,"不论你面向哪边…",sz=20,b=True,c=DARK)
# Big highlight bar
hi=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.15),Inches(2.85),Inches(4.30),Inches(0.85))
hi.fill.solid();hi.fill.fore_color.rgb=COMPRED;hi.line.fill.background()
tb(s,5.15,2.92,4.30,0.7,"🔴 红针都指北!",sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,5.15,3.85,4.4,0.35,"Red needle ALWAYS points NORTH!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,5.15,4.20,4.4,0.35,"⭐ 这就是它神奇的地方!",sz=13,b=True,c=NAVY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.80,
    "你们发现了什么?",
    "全班喊 3 次: 红针 — 永远 — 指北!")
n+=1;pn(s,n)
notes(s,"小结 (3 分钟):\n• 让学生自己总结刚才的发现\n• 老师小结: 「不论我们站在哪里, 面向什么方向, 红色磁针都会指向同一个方向 — 北方!」\n• 介绍指南针 2 个组成: 方位盘 (圆盘) + 磁针 (动的针)\n• 红色 = 北 (国际通用)\n• 这就是为什么叫 「指南针」 — 但其实它是 「指北针」 哦!")

# 15. COMPASS HISTORY VIDEO (厉害的指南针)
s=ns();bg(s,NAVY)
tb(s,1,0.45,8,0.6,"🎬 看视频  Watch a Video",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.05,8,0.4,"《厉害的指南针》  The Amazing Compass — 古代到现代",sz=15,b=True,c=GOLD,a=PP_ALIGN.CENTER)
# Video frame
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(1.65),Inches(7.0),Inches(2.3))
sh.fill.solid();sh.fill.fore_color.rgb=DARK;sh.line.color.rgb=WHITE;sh.line.width=Pt(3)
tb(s,1.5,2.30,7.0,1.0,"▶️",sz=100,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.5,3.55,7.0,0.35,"🔗 D3 resource/厉害的指南针_科普_大班_视频.mp4",sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
# Pre-watch questions
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.10),Inches(9.0),Inches(1.0))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=GOLD;sh2.line.width=Pt(2)
tb(s,0.65,4.18,8.7,0.3,"👂 听 / 看 4 个问题:",sz=12,b=True,c=NAVY)
tb(s,0.65,4.48,4.5,0.3,"1. 最早的指南针叫什么?",sz=11,c=DARK)
tb(s,5.00,4.48,4.5,0.3,"2. 后来又发明了什么?",sz=11,c=DARK)
tb(s,0.65,4.78,4.5,0.3,"3. 磁铁针 + ___ = 指南针?",sz=11,c=DARK)
tb(s,5.00,4.78,4.5,0.3,"4. 人们在哪里用指南针?",sz=11,c=DARK)
n+=1;pn(s,n)
notes(s,"看视频 (3-4 分钟):\n• 视频文件: D3 resource/厉害的指南针_科普_大班_视频.mp4\n• 看视频前布置 4 个问题让学生注意\n• 视频内容: 指南针的演变 (司南 → 指南鱼 → 磁铁针 → 指南针)\n• 看完后用下一页 review 答案")

# 16. COMPASS HISTORY — 4 stages reveal
s=ns();bg(s,CREAM);hb(s,"📜 指南针的故事  The Compass Story",COMPRED)
tb(s,0.4,0.85,9.2,0.4,"中国古代四大发明之一!  One of China's 4 Great Inventions!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# 4 stages timeline
stages=[("1️⃣","司南","Sinan","勺子 + 盘\n容易破!",BROWN),
        ("2️⃣","指南鱼","Compass Fish","磁铁片\n做的小鱼",E_CL),
        ("3️⃣","磁铁针","Magnetic Needle","只有一根\n磁针",NAVY),
        ("4️⃣","指南针","Compass","磁针 + 方位盘\n现在用的!",COMPRED)]
for i,(num,cn,en,desc,cl) in enumerate(stages):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.40),Inches(2.20),Inches(3.15))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    # Compact number marker (top-left corner)
    pill(s,x+0.10,1.45,0.55,0.30,num.replace("️⃣","").strip(),cl,sz=12)
    # LARGER image area
    ib(s,x+0.10,1.83,2.00,1.85,"📷")
    # Title + EN on one line
    tb(s,x+0.10,3.75,2.0,0.40,f"{cn}",sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.10,2.0,0.25,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    # Single-line desc (joined)
    desc_one = " · ".join(desc.split("\n"))
    tb(s,x+0.10,4.32,2.0,0.20,desc_one,sz=9,c=DARK,a=PP_ALIGN.CENTER)
# Arrows between
for ax in [2.55,4.90,7.25]:
    tb(s,ax,2.55,0.20,0.5,"→",sz=22,b=True,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.55,
    "古人真聪明! 你最喜欢哪个?",
    "我喜欢 ___ , 因为 ___ 。")
n+=1;pn(s,n)
notes(s,"指南针的演变 (3 分钟):\n• 司南: 最早, 像勺子放在盘上, 但容易破碎\n• 指南鱼: 用磁铁片做成鱼形, 放水里浮着指北\n• 磁铁针: 只用一根磁针 (更方便)\n• 指南针: 磁针 + 方位盘 = 现代指南针\n• 这是中国古代四大发明 (造纸/火药/指南针/印刷)\n• 让学生说: 我喜欢哪一个? 为什么?")

# 17. WHEN DO PEOPLE USE COMPASS?
s=ns();bg(s,CREAM);hb(s,"🌍 在哪里用指南针?  Where Do People Use It?",COMPRED)
tb(s,0.4,0.85,9.2,0.4,"指南针在这些地方都很重要!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Compass is important in many places!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
uses=[("⛺","在野外","In the wild",PINE),
      ("🏜️","在沙漠","In the desert",SUN),
      ("⚔️","在战场","On the battlefield",ALERT),
      ("⛵","在海上","At sea",NAVY)]
for i,(em,cn,en,cl) in enumerate(uses):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(2.20),Inches(2.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.85,2.1,1.2,em,sz=72,a=PP_ALIGN.CENTER)
    ib(s,x+0.15,3.10,1.90,0.85,"📷")
    tb(s,x+0.1,4.00,2.0,0.4,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.55,
    "你想在哪里用指南针?",
    "我想在 ___ 用指南针, 因为 ___ 。")
n+=1;pn(s,n)
notes(s,"应用 (2-3 分钟):\n• 让学生看 4 个场景\n• 提问: 「为什么这些地方需要指南针?」\n• 关键: 这些地方没有路标 / GPS, 容易迷路\n• 句型: 我想在 ___ 用指南针")

# 18. COMPASS — how to use (3 simple steps)
s=ns();bg(s,CREAM);hb(s,"📍 怎么用指南针?  How to Use",COMPRED)
tb(s,0.4,0.85,9.2,0.4,"3 步使用法  Just 3 Steps",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
steps=[("1️⃣","平拿","Hold flat","平拿指南针 ✋\nHold flat",NAVY),
       ("2️⃣","等红针","Wait","等红针停下来\nWait for red",COMPRED),
       ("3️⃣","跟着走","Follow","红针指 = 北 → 跟着走!\nRed = North, follow!",E_CL)]
for i,(num,cn,en,desc,cl) in enumerate(steps):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(3.0),Inches(2.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.55,2.8,0.7,num,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.30,2.8,0.5,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.85,2.8,0.4,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    lns=desc.split('\n')
    tb(s,x+0.15,3.30,2.8,0.4,lns[0],sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    if len(lns)>1: tb(s,x+0.15,3.70,2.8,0.4,lns[1],sz=10,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,
    "我们一起做 — 平拿! 等! 走!",
    "假装手里有指南针 — 跟着老师做 3 步!")
n+=1;pn(s,n)

# 15. SESSION 2 DIVIDER
s=div("Session 2  下午 45 min","🔄 复习 + 我会认 5 / 我会写 4",SUN,"📖");n+=1;pn(s,n)

# 16. REVIEW — Baamboozle game (teacher inserts link)
s=ns();bg(s,CREAM);hb(s,"🎮 复习游戏  Review · Baamboozle",SUN)
# Big centered placeholder for teacher's link / screenshot
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(1.30),Inches(7.0),Inches(3.3))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
tb(s,1.5,1.85,7.0,0.6,"🎲",sz=80,a=PP_ALIGN.CENTER)
tb(s,1.5,2.85,7.0,0.5,"老师在这里放 Baamboozle 链接",sz=20,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,1.5,3.40,7.0,0.4,"Teacher: paste your Baamboozle link here",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1.5,3.90,7.0,0.4,"🔗 https://www.baamboozle.com/...",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.80,
    "我们来玩复习游戏!",
    "分组抢答 — 答对得分!")
n+=1;pn(s,n)
notes(s,"复习游戏 (10-15 分钟):\n• 老师提前在 Baamboozle 上创建一个复习游戏 (或用已有的)\n• 题目涵盖: 4 个方向 / 太阳从哪边升起 / 红针指哪 / 指南针组成 / 司南是什么 / ...\n• 把 Baamboozle 链接粘贴到这一页\n• 上课时直接打开链接玩\n• 全班分 2-3 组, 抢答得分")

# 17. POINT-TO-WORD GAME
s=ns();bg(s,CREAM);hb(s,"🎯 指对字!  Point to the Right Word!",E_CL)
tb(s,0.4,0.85,9.2,0.4,"老师说英文, 你指中文字!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Teacher says English — you point to the Chinese character!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
words=[("东","E·East",E_CL),("南","S·South",S_CL),("西","W·West",W_CL),("北","N·North",N_CL),("指南针","Compass",COMPRED)]
for i,(w,en,cl) in enumerate(words):
    x=0.4+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(1.80),Inches(2.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    sz=72 if len(w)==1 else 32
    y_off=2.05 if len(w)==1 else 2.35
    tb(s,x+0.05,y_off,1.7,1.3,w,sz=sz,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,1.7,0.5,en,sz=14,b=True,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.55,
    "「North」在哪? 「Compass」在哪?",
    "用手指 — 大声读出来!")
n+=1;pn(s,n)
notes(s,"游戏 (3 分钟):\n• 老师喊英文, 学生指对应的中文字 + 大声读\n• 反过来: 老师喊中文, 学生说英文\n• 速度慢的小组先练习, 速度快的小组挑战")

# 18-22. WORD CARDS — 我会认 (5)
read_data=[
    ("东","dōng","East",       "🌅","早上, 太阳从东边升起。"),
    ("南","nán", "South",      "🔥","中午, 太阳在南边 (北半球)。"),
    ("西","xī",  "West",       "🌇","下午, 太阳从西边落下。"),
    ("北","běi", "North",      "❄️","北边比较冷, 指南针红针指北。"),
    ("指南针","zhǐ nán zhēn","Compass","🧭","指南针的红针指向北。"),
]
for w,py,en,em,sent in read_data:
    s=word_card_read(w,py,en,em,sent,color=NAVY);n+=1;pn(s,n)

# 23-26. WORD CARDS — 我会写 (4 — 东南西北 only)
write_data=[
    ("东","dōng","East", 5,"5 笔: 横 ㇆ 撇点点  ·  Two dots like sun rays"),
    ("南","nán","South", 9,"9 笔: 上面十 + 下面是 「冂」+「𢆉」  ·  Like a tent shape"),
    ("西","xī", "West",  6,"6 笔: 像一个杯子 — 下面一横 + 凵 + 儿  ·  Like a cup"),
    ("北","běi","North", 5,"5 笔: 两个人背靠背 ⇆  ·  Two people back-to-back"),
]
for w,py,en,sc,hint in write_data:
    s=word_card_write(w,py,en,sc,hint,color=NAVY);n+=1;pn(s,n)

# 27. CHRISTMAS COMPASS — activity launch
s=ns();bg(s,CREAM);hb(s,"🎁 圣诞老人加工厂方向指南针  Santa's Workshop Compass",COMPRED)
# LEFT: full worksheet image
img(s,0.3,0.95,4.6,4.3,os.path.join(BASE,"pics","santa_compass-1.png"),"📷 Santa workshop grid worksheet")
# RIGHT: how-to-play card
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(0.95),Inches(4.6),Inches(0.45))
sh.fill.solid();sh.fill.fore_color.rgb=COMPRED;sh.line.fill.background()
tb(s,5.2,1.00,4.4,0.4,"🎯 怎么玩  How to Play",sz=15,b=True,c=WHITE)
tf=tb(s,5.2,1.55,4.4,0.4,"1️⃣ 找到「开始」格 ⭕  Find START",sz=13,b=True,c=DARK)
ap(tf,"2️⃣ 听 / 读方向指令  Read direction",sz=13,b=True,c=DARK)
ap(tf,"3️⃣ 一格一格走  Move 1 grid at a time",sz=13,b=True,c=DARK)
ap(tf,"4️⃣ 说: 我到了 ___ !  I arrived at ___",sz=13,b=True,c=COMPRED)
# Cardinal-direction practice questions (K-5 friendly)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(3.10),Inches(4.6),Inches(0.45))
sh2.fill.solid();sh2.fill.fore_color.rgb=NAVY;sh2.line.fill.background()
tb(s,5.2,3.15,4.4,0.4,"🧭 试一试 (基础 — 东南西北)",sz=14,b=True,c=WHITE)
tf2=tb(s,5.2,3.65,4.4,0.4,"Q1. 从开始, 向北走 4 格 → ?",sz=12,b=True,c=DARK)
ap(tf2,"Q2. 向南走 2 格 → ?",sz=12,b=True,c=DARK)
ap(tf2,"Q3. 向西走 4 格 → ?",sz=12,b=True,c=DARK)
ap(tf2,"💬 我向 ___ 走 ___ 格, 我到了 ___ 。",sz=11,b=True,c=NAVY)
n+=1;pn(s,n)
notes(s,"圣诞工厂方向活动 (8-10 分钟):\n• 老师打印 worksheet 每位学生 1 张 (来源: twinkl 二年级数学指南针活动)\n• 先一起做 Q1: 「从开始向北走 4 格」 — 全班数 1-2-3-4\n• 答案是弹球 (pinballs) 那一格\n• Q2 向南走 2 格 → 滑板车; Q3 向西走 4 格 → 遥控汽车\n• 让学生小组互相检查\n• 句型: 「我向北走 4 格, 我到了弹球!」\n• 进阶版本在下一页 (东北/西南等 8 个方向)")

# 28. CHRISTMAS COMPASS — answer reveal + 8 directions extension
s=ns();bg(s,CREAM);hb(s,"🎁 进阶挑战 + 答案  Challenge + Answers",COMPRED)
# LEFT: 8-direction grid (3x3 cells, no compass rose to avoid overlap)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(3.55))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=COMPRED;sh.line.width=Pt(2.5)
tb(s,0.4,1.05,4.2,0.4,"🧭 8 个方向  8 Directions",sz=15,b=True,c=COMPRED,a=PP_ALIGN.CENTER)
# 3x3 grid of direction cells
cells=[("西北","NW",W_CL),("北","N",N_CL),("东北","NE",N_CL),
       ("西","W",W_CL),  ("🧭","",GRAY),("东","E",E_CL),
       ("西南","SW",W_CL),("南","S",S_CL),("东南","SE",E_CL)]
gx=0.6;gy=1.55;cw=1.27;ch=0.85
for idx,(cn,en,cl) in enumerate(cells):
    r=idx//3;c=idx%3
    x=gx+c*cw;y=gy+r*ch
    sq=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(cw-0.05),Inches(ch-0.05))
    sq.fill.solid();sq.fill.fore_color.rgb=WHITE if cn!="🧭" else GOLD
    sq.line.color.rgb=cl;sq.line.width=Pt(1.5)
    if cn=="🧭":
        tb(s,x,y+0.18,cw-0.05,0.5,cn,sz=28,a=PP_ALIGN.CENTER)
    else:
        tb(s,x,y+0.10,cw-0.05,0.4,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
        tb(s,x,y+0.48,cw-0.05,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,4.2,0.3,"💡 北+东=东北 · 南+西=西南",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
# RIGHT: 7 questions with answers (compact)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.40))
sh2.fill.solid();sh2.fill.fore_color.rgb=NAVY;sh2.line.fill.background()
tb(s,5.0,0.99,4.6,0.35,"💡 答案  Answers",sz=13,b=True,c=WHITE)
qa=[
    ("Q1","北 4 格","→ 弹球",E_CL),
    ("Q2","东北 1 格","→ 多米诺骨牌",N_CL),
    ("Q3","南 2 格","→ 滑板车",S_CL),
    ("Q4","西 4 格","→ 遥控汽车",W_CL),
    ("Q5","东南 2 格","→ 木制小汽车",E_CL),
    ("Q6","滑板 → 泰迪熊","西南 1 格",W_CL),
    ("Q7","泰迪熊 → 多米诺","东北 4 格",N_CL),
]
for i,(q,d,a,cl) in enumerate(qa):
    y=1.45+i*0.43
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(y),Inches(4.8),Inches(0.38))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(1.5)
    pill(s,5.0,y+0.04,0.55,0.30,q,cl,sz=11)
    tb(s,5.65,y+0.05,1.85,0.28,d,sz=11,b=True,c=DARK)
    tb(s,7.55,y+0.05,2.1,0.28,a,sz=11,b=True,c=cl)
teacher_student_bar(s,4.65,
    "你最喜欢哪个玩具? 走过去!",
    "「我向 ___ 走 ___ 格, 我到了 ___ 。」")
n+=1;pn(s,n)
notes(s,"答案 + 进阶 (5-8 分钟):\n• 介绍 4 个对角方向: 东北/东南/西北/西南 (8 个方向指南针)\n• 公式: 北+东=东北, 北+西=西北, 南+东=东南, 南+西=西南\n• 让学生自己出题考别人 (Q8 open-ended): 写一组方向指令, 同桌猜终点\n• 收集 worksheet 检查\n• 句型练习: 我向 ___ 走 ___ 格 — 这正好是 Project 1 (Treasure Hunt) 的预热!")

# 27→ 29 (renumbered). SESSION 3 DIVIDER
s=div("Session 3  下午 60-90 min","🛠️ 3 个项目  Treasure · Map · Compass",BROWN,"🎒");n+=1;pn(s,n)

# 28. PROJECTS OVERVIEW
s=ns();bg(s,CREAM);hb(s,"🛠️ 动手时间！  Hands-On Time — 3 Projects",BROWN)
projects=[
    ("PROJECT 1","🎯 寻宝任务","Treasure Hunt","按指令找位置\nFollow directions",WARM,COMPRED,"基础 / Basic"),
    ("PROJECT 2","🏙️ 设计我的城市","Design My City","剪 + 贴 + 标方向\nCut, glue, label",RGBColor(0xFF,0xE0,0xB2),NAVY,"中级 / Mid"),
    ("PROJECT 3","🧲 自制指南针","DIY Compass","针 + 磁铁 + 水\nNeedle + magnet + water",RGBColor(0xDC,0xED,0xC8),OK,"进阶 / Advanced"),
]
for i,(lbl,nm,en,d,bgc,cl,lvl) in enumerate(projects):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(3.1),Inches(4.15))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,2.9,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.4,2.9,0.6,nm,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.0,2.9,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.85,2.45,1.4,0.35,lvl,cl,sz=10)
    ib(s,x+0.2,2.95,2.8,1.2,"📷 示范")
    ls=d.split('\n')
    tf=tb(s,x+0.15,4.20,2.85,0.4,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for ln in ls[1:]:ap(tf,ln,sz=12,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# 29. PROJECT 1 — Treasure Hunt
s=ns();bg(s,CREAM);hb(s,"🎯 Project 1: 寻宝任务  Treasure Hunt",COMPRED)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=COMPRED;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.1,"🃏 指令卡片  Direction cards",sz=13,c=DARK)
ap(tf,"🎁 小奖品  Hidden treasures",sz=13,c=DARK)
ap(tf,"🧭 指南针 (Project 3 做的)",sz=13,c=DARK)
ap(tf,"📏 「3 步」「5 步」步数卡",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 玩法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.1,"1️⃣ 老师藏 4-5 个小物品",sz=13,c=DARK)
ap(tf2,"2️⃣ 给学生指令卡:",sz=13,c=DARK)
ap(tf2,"     「向北 5 步, 向东 3 步」",sz=12,c=NAVY)
ap(tf2,"3️⃣ 学生用指南针走过去!",sz=13,c=DARK)
ap(tf2,"4️⃣ 找到 — 大声说指令",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.35))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=COMPRED;sh3.line.width=Pt(2)
tb(s,0.5,3.95,9,0.35,"🗣️ 展示句型  Say These:",sz=14,b=True,c=COMPRED)
tb(s,0.5,4.35,4.5,0.35,"·  我向 ___ 走 ___ 步。",sz=14,c=DARK)
tb(s,0.5,4.65,4.5,0.35,"·  我找到了 ___ !",sz=14,c=DARK)
tb(s,5.2,4.35,4.5,0.35,"·  这是 ___ 边 (东/南/西/北)。",sz=14,c=DARK)
tb(s,5.2,4.65,4.5,0.35,"·  Treasure is in the ___ corner.",sz=14,c=DARK)
n+=1;pn(s,n)
notes(s,"老师准备 (低 prep):\n• 提前藏 4-5 件小物品 (橡皮、贴纸、糖果)\n• 写 4-5 张指令卡: 「向北 5 步, 向东 3 步」\n• 学生用 Project 3 做的指南针走\n• 找到的人大声读自己的指令")

# 30. PROJECT 2 — Design a City (with directions)
s=ns();bg(s,CREAM);hb(s,"🏙️ Project 2: 设计我的城市  Design My City",NAVY)
# LEFT: city map preview image
img(s,0.3,0.95,4.4,3.40,os.path.join(BASE,"pics","design_city_map-3.png"),"📷 City map grid w/ N E S W compass")
# RIGHT TOP: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(0.40))
sh.fill.solid();sh.fill.fore_color.rgb=NAVY;sh.line.fill.background()
tb(s,4.95,0.99,4.65,0.35,"🧺 材料  Materials",sz=13,b=True,c=WHITE)
tf=tb(s,4.95,1.45,4.7,1.6,"📄 城市地图 (印好) + 建筑卡片",sz=12,c=DARK)
ap(tf,"   City map + cut-out building cards",sz=10,c=GRAY)
ap(tf,"✂️ 剪刀  Scissors",sz=12,c=DARK)
ap(tf,"📎 胶水 / 双面胶  Glue / tape",sz=12,c=DARK)
ap(tf,"🖍️ 彩笔 (写方向)  Markers (label)",sz=12,c=DARK)
# RIGHT MIDDLE: steps
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(3.10),Inches(4.85),Inches(0.40))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,4.95,3.14,4.65,0.35,"👉 做法  Steps",sz=13,b=True,c=WHITE)
tf2=tb(s,4.95,3.55,4.7,0.4,"1️⃣ 看地图 — 找 N E S W (上北下南左西右东)",sz=11,c=DARK)
ap(tf2,"2️⃣ 剪喜欢的建筑 (学校 / 医院 / 商店…)",sz=11,c=DARK)
ap(tf2,"3️⃣ 贴在地图 — 想想它在哪个方向",sz=11,c=DARK)
ap(tf2,"4️⃣ 用中文写出每栋楼的方向",sz=11,b=True,c=NAVY)
# BOTTOM: speaking patterns (full width)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.45),Inches(9.4),Inches(0.75))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=NAVY;sh3.line.width=Pt(2)
tb(s,0.45,4.50,9.1,0.30,"🗣️ 展示句型  Say These:",sz=12,b=True,c=NAVY)
tb(s,0.45,4.80,4.5,0.30,"·  学校在城市的 ___ 边 。",sz=12,c=DARK)
tb(s,0.45,5.05,4.5,0.30,"·  医院在学校的 ___ 边 。",sz=12,c=DARK)
tb(s,5.15,4.80,4.5,0.30,"·  这是我设计的城市!",sz=12,c=DARK)
tb(s,5.15,5.05,4.5,0.30,"·  My city has a ___ in the ___ .",sz=12,c=DARK)
n+=1;pn(s,n)
notes(s,"中等 prep (15-20 分钟):\n• 老师提前打印:\n  - 每位学生 1 张地图 (page 3 of 'Design a City' PDF)\n  - 每位学生 1 张建筑卡片 (page 4 — 剪开)\n  - 来源: Twinkl 'Design a City Activity for K-2nd Grade'\n• 流程:\n  1. 先一起看城市地图: 找到 N/E/S/W 4 个方向\n  2. 学生选 5-8 个建筑剪下来\n  3. 想一想这个城市需要什么? 学校放哪? 医院放哪?\n  4. 贴在地图上 — 注意方向\n  5. 在每个建筑旁边用中文标方向 (例: 学校在西边)\n• 完成后 share: 「我的城市 ___ 在 ___ 边」")

# 31. PROJECT 3 — Make a Compass
s=ns();bg(s,CREAM);hb(s,"🧲 Project 3: 自制指南针  DIY Compass",OK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=OK;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.1,"🪡 缝衣针  Sewing needle",sz=13,c=DARK)
ap(tf,"🧲 磁铁  Magnet (老师准备)",sz=13,c=DARK)
ap(tf,"🟫 软木 / 泡沫小片  Cork / foam",sz=13,c=DARK)
ap(tf,"🥣 一碗水  Bowl of water",sz=13,c=DARK)
ap(tf,"⚠️ 老师监督 — 针很尖!",sz=12,b=True,c=ALERT)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法 4 步  4 Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.2,"1️⃣ 用磁铁朝同一方向摩擦针 30 次",sz=13,c=DARK)
ap(tf2,"     Rub same direction, 30 times",sz=10,c=GRAY)
ap(tf2,"2️⃣ 把针放软木 / 泡沫上",sz=13,c=DARK)
ap(tf2,"3️⃣ 把软木放水面上",sz=13,c=DARK)
ap(tf2,"4️⃣ 等针稳定 — 它指向北! 🧭",sz=13,b=True,c=COMPRED)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.95),Inches(9.4),Inches(1.25))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=OK;sh3.line.width=Pt(2)
tb(s,0.5,4.05,9,0.35,"🎬 视频教程  Video:",sz=14,b=True,c=OK)
tb(s,0.5,4.40,9,0.35,"🔗 https://www.youtube.com/shorts/YZiVvUFkibE",sz=12,b=True,c=NAVY)
tb(s,0.5,4.75,9,0.35,"🗣️ 我做的指南针指向 ___ 。 / My compass points to ___.",sz=12,b=True,c=DARK)
n+=1;pn(s,n)
notes(s,"老师准备 (中等 prep):\n• 每组 1 块磁铁, 5-10 根针, 1 个浅碗水, 软木\n• 安全: 老师全程监督, 摩擦后老师收回\n• 视频提前看一遍\n• 失败常见原因: 摩擦次数太少 / 方向不一致 / 软木太大不平")

# 32. CLOSING — mission complete
s=ns();bg(s,NAVY)
tb(s,1,0.6,8,0.7,"🏆 任务完成!",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.35,8,0.5,"Mission Complete!",sz=20,c=WARM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.0),Inches(3.0),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.color.rgb=WHITE;sh.line.width=Pt(4)
tb(s,3.5,2.4,3.0,0.6,"🧭",sz=80,a=PP_ALIGN.CENTER)
tb(s,3.5,3.6,3.0,0.5,"方向大师",sz=22,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,3.5,4.1,3.0,0.4,"Direction Master",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"明天: Day 4 — 食物与水 / Day 4: Food & Water",sz=14,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# === Save ===
out=os.path.join(BASE,"day3_compass.pptx")
prs.save(out)
print(f"Saved {out}  ({n} slides)")
