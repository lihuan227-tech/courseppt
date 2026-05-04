#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
野外生存与探险 Day 4 — 紧急情况怎么办?  (Emergencies: What to Do?)
K-5 immersive lesson.
Session 1: Weather dangers · Wild animals · Don't touch/eat · Lost
Session 2: Review + 我会认 (天气/危险/迷路/食物) + 我会写 (天气/危险)
Session 3: Safety skits + Safety booklet + Bonus activities
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Safety/Emergency palette ---
PINE   = RGBColor(0x1E,0x4D,0x2B)
SUN    = RGBColor(0xE0,0x7A,0x2C)
CREAM  = RGBColor(0xFD,0xF6,0xE3)
BROWN  = RGBColor(0x6B,0x44,0x23)
SKY    = RGBColor(0x4A,0x90,0xD9)
SUNYEL = RGBColor(0xF5,0xC2,0x42)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
WARM   = RGBColor(0xFF,0xF3,0xE0)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
NAVY   = RGBColor(0x1A,0x23,0x7E)
# Emergency-themed
DANGER  = RGBColor(0xC6,0x28,0x28)   # warning red (dominant)
CAUTION = RGBColor(0xF9,0xA8,0x25)   # caution amber
SAFE    = RGBColor(0x38,0x8E,0x3C)   # safe green
CALM    = RGBColor(0x1E,0x88,0xE5)   # calm blue
ALERT   = DANGER

BASE = "/Users/Huan/0 projects/summercourse/Chinese/野外生存与探险wilderness_pbl"

# === Helpers (same vocab as day3) ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def img(s,l,t,w,h,path,fallback="📷"):
    if os.path.exists(path):
        s.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(h))
    else:
        ib(s,l,t,w,h,fallback)
def hb(s,txt,c=DANGER,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,txt):
    s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.8,9,0.8,sub,sz=20,c=WHITE,a=PP_ALIGN.CENTER);return s
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)

def teacher_student_bar(s,t,teacher_q,student_action,color=DANGER):
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=color;sf.line.width=Pt(2)
    tb(s,0.45,t+0.04,4.5,0.25,"👩‍🏫 老师问 Teacher asks:",sz=10,b=True,c=color)
    tb(s,0.45,t+0.27,4.5,0.28,teacher_q,sz=12,b=True,c=DARK)
    tb(s,5.0,t+0.04,4.6,0.25,"🧒 学生 Student does:",sz=10,b=True,c=SUN)
    tb(s,5.0,t+0.27,4.6,0.28,student_action,sz=12,b=True,c=DARK)

def danger_safe_slide(title_cn,title_en,subject_em,subject_cn,subject_en,
                     dangers,safe_actions,img_hint,header_color,brainstorm=False):
    """One slide per topic — image LEFT, danger card + safe actions card RIGHT.
    brainstorm=True: hide dangers/actions, show ❓ prompts for student brainstorm.
    """
    s=ns();bg(s,CREAM);hb(s,f"{subject_em} {title_cn}  {title_en}",header_color)
    ib(s,0.4,1.0,4.4,3.6,f"📷 {img_hint}")
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
    tb(s,5.15,1.10,4.4,0.35,"⚠️ 危险?  What's the danger?",sz=13,b=True,c=DANGER)
    if brainstorm:
        tb(s,5.15,1.55,4.4,1.1,"❓",sz=72,c=DANGER,a=PP_ALIGN.CENTER)
    else:
        y=1.50
        for d in dangers:
            tb(s,5.15,y,4.4,0.35,f"·  {d}",sz=12,b=True,c=DARK);y+=0.32
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(2.85),Inches(4.6),Inches(1.75))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SAFE;sh2.line.width=Pt(2.5)
    tb(s,5.15,2.95,4.4,0.35,"✅ 怎么办?  What to do?",sz=13,b=True,c=SAFE)
    if brainstorm:
        tb(s,5.15,3.35,4.4,1.1,"❓",sz=72,c=SAFE,a=PP_ALIGN.CENTER)
    else:
        y=3.35
        for a in safe_actions:
            tb(s,5.15,y,4.4,0.35,f"·  {a}",sz=12,b=True,c=DARK);y+=0.32
    if brainstorm:
        teacher_student_bar(s,4.75,
            f"看到 {subject_cn} — 你觉得?",
            "Think-Pair-Share: 危险是 ___, 我会 ___ 。")
    else:
        teacher_student_bar(s,4.75,
            f"看到 {subject_cn} 怎么办?",
            f"我看到 {subject_cn}, 我不 ___, 我会 ___ 。")
    return s

def word_card_read(w,py,en,img_emoji,sent,color=DANGER):
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(2.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,0.5,1.3,4.2,2.0,img_emoji,sz=130,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(2.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=color;sh2.line.width=Pt(2.5)
    char_sz = 110 if len(w)==1 else (66 if len(w)==2 else 54)
    tb(s,5.1,1.15,4.4,1.5,w,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,5.1,2.70,4.4,0.4,f"{py}  ·  {en}",sz=18,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,5.1,3.15,4.4,0.4,"👉 跟我读 3 次！",sz=14,c=color,a=PP_ALIGN.CENTER)
    sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.3))
    sh3.fill.solid();sh3.fill.fore_color.rgb=WHITE;sh3.line.color.rgb=color;sh3.line.width=Pt(2)
    tb(s,0.6,3.95,2.0,0.4,"💬 例句 Example",sz=14,b=True,c=color)
    tb(s,0.6,4.35,8.8,0.6,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,strokes_count,strokes_hint,color=DANGER):
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    char_sz = 160 if len(w)==1 else 100
    tb(s,0.5,1.1,4.2,2.4,w,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.55,4.2,0.5,f"{py}  ·  {en}",sz=22,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.10,4.2,0.4,f"{strokes_count} 笔 / {strokes_count} strokes",sz=16,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.50,4.2,0.4,strokes_hint,sz=11,c=DARK,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,5.15,1.1,4.4,0.4,"3 步练习  3 Steps",sz=16,b=True,c=color)
    tb(s,5.15,1.5,4.4,0.4,"1️⃣ 看老师写  Watch teacher",sz=13,c=DARK)
    tb(s,5.15,1.85,4.4,0.4,"2️⃣ 用手指空中写  Air-write",sz=13,c=DARK)
    tb(s,5.15,2.20,4.4,0.4,"3️⃣ 在田字格写 3 次",sz=13,c=DARK)
    for i in range(4):
        x=5.0+i*1.15
        sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.85),Inches(1.05),Inches(1.05))
        sq.fill.solid();sq.fill.fore_color.rgb=WHITE;sq.line.color.rgb=color;sq.line.width=Pt(1.5)
        ln1=s.shapes.add_connector(1,Inches(x),Inches(2.85+0.525),Inches(x+1.05),Inches(2.85+0.525))
        ln1.line.color.rgb=LGRAY;ln1.line.width=Pt(0.5);ln1.line.dash_style=2
        ln2=s.shapes.add_connector(1,Inches(x+0.525),Inches(2.85),Inches(x+0.525),Inches(2.85+1.05))
        ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.5);ln2.line.dash_style=2
    tb(s,5.0,4.0,4.6,0.3,"在田字格里写 3 次 ↓",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    teacher_student_bar(s,4.45,
        f"和我一起写「{w}」",
        "看 → 空中写 → 田字格写 3 次")
    return s

# ========================================================================
#                              SLIDES
# ========================================================================
n=0

# 1. COVER
s=ns();bg(s,DANGER)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.4,8,0.5,"DAY 4",sz=18,b=True,c=CAUTION,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"⚠️ 紧急情况怎么办?",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"What to Do in Emergencies",sz=20,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"🆘 安全探险家  Safety Explorer",sz=24,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"天气 · 动物 · 别碰别吃 · 迷路怎么办",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"Weather · Animals · Don't Touch · Lost?",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"野外生存与探险 · Wilderness Survival",sz=14,b=True,c=CAUTION,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「探险家们! 在野外, 总会遇到一些紧急情况。今天我们要学会 4 件事 — 怎么保护自己!」\n• 用提问 hook: 「你害怕过什么? 在野外有什么危险?」")

# 2. SESSION 1 LEARNING GOALS
s=ns();bg(s,CREAM);hb(s,"🎯 Session 1 学习目标  Today's Goals",DANGER)
tb(s,0.4,0.85,9.2,0.4,"上完 Session 1, 你可以…",sz=18,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"After Session 1, you can…",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
goals=[
    ("1️⃣","认识天气和危险","Know weather + dangers (sun/rain/wind/heat/cold)",CAUTION),
    ("2️⃣","知道怎么应对野生动物","Know what to do with wild animals (snake/bear/bee…)",PINE),
    ("3️⃣","知道不能碰 / 吃的东西","Know what NOT to touch / eat",DANGER),
    ("4️⃣","学会迷路怎么办","Know what to do if you get lost",CALM),
]
for i,(num,cn,en,cl) in enumerate(goals):
    y=1.70+i*0.78
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.70))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.55,y+0.13,0.7,0.4,num,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.05,3.5,0.35,cn,sz=17,b=True,c=cl)
    tb(s,5.00,y+0.20,4.5,0.35,en,sz=12,c=DARK)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 跟读 4 个目标\n• 「这 4 件事 — 都是真的紧急情况! 学会了, 你就是安全探险家!」")

# 3. SESSION 1 DIVIDER
s=div("Session 1  上午","⚠️ 4 个紧急情况  Weather · Animals · Don't Touch · Lost",DANGER,"📖");n+=1;pn(s,n)

# ===== TOPIC 1: WEATHER =====
# 4. WEATHER intro — 4 weather types (predict)
s=ns();bg(s,CREAM);hb(s,"🌤️ 天气会带来什么危险?  Weather Dangers",CAUTION)
tb(s,0.4,0.85,9.2,0.4,"在野外的时候, 如果遇到这样的天气, 可能会有哪些危险?",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"What dangers can these weathers bring?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
weathers=[("☀️","晴天 / 很热","Sunny / Hot",SUN),
          ("🌧️","下雨","Rainy",CALM),
          ("💨","大风","Windy",GRAY),
          ("🥶","很冷","Cold",CALM)]
for i,(em,cn,en,cl) in enumerate(weathers):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.20),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.95,2.1,1.0,em,sz=64,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,2.1,0.5,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,2.1,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,
    "你觉得哪个最危险?",
    "Think-Pair-Share: 我觉得 ___ 最危险, 因为 ___ 。")
n+=1;pn(s,n)
notes(s,"hook (3 分钟):\n• 不要立刻给答案\n• 让学生小组讨论, 收集所有想法\n• 「我们一个一个看 — 看看你猜的对不对!」")

# --- WEATHER (brainstorm → reveal pair for each) ---
weathers_data=[
    # (title, en, em, dangers, safe_actions, img_hint, header_color, notes_extra)
    ("晴天 / 很热","Sunny / Hot","☀️",
        ["晒伤","中暑","口渴"],
        ["戴帽子 + 防晒霜","多喝水","在树下休息"],
        "孩子在大太阳下出汗",CAUTION,
        "中暑表现: 头晕, 想吐, 没力气 — 立刻找大人"),
    ("下雨","Rain","🌧️",
        ["地滑, 容易摔倒","衣服湿 → 容易着凉","看不清路"],
        ["带雨衣 / 雨伞","小心走路, 不跑","找干燥的地方躲"],
        "下雨天的森林步道",CALM,
        "不要玩水洼; 大雨要小心山洪"),
    ("⚡ 打雷闪电","Thunder & Lightning","⚡",
        ["雷会击中高的东西","金属 (伞柄) 会引雷","树下 / 水边都危险"],
        ["不站树下! 不站水边!","蹲下抱头 (脚并拢)","远离金属 (放下伞)","快进屋子 / 帐篷"],
        "雷电闪电的天空",DANGER,
        "「30 / 30 规则」: 看到闪电 → 30秒内听到雷 = 危险; 雷停后等 30 分钟才出来"),
    ("大风","Windy","💨",
        ["树枝可能掉下来","帐篷可能被吹动","沙子进眼睛"],
        ["远离大树 / 建筑","固定帐篷 (用石头压)","闭眼 / 戴太阳镜"],
        "大风吹动的树和帐篷",GRAY,
        "可以联系 Day 2 的帐篷搭建技巧"),
    ("很冷","Cold","🥶",
        ["着凉 / 感冒","手脚发抖","容易冻伤 (耳朵 / 手指)"],
        ["多穿衣服 (洋葱式)","戴帽子 + 手套","活动一下暖暖身体"],
        "雪地里穿厚衣服的孩子",CALM,
        "洋葱式穿衣 = 一层一层 — 热了脱一件"),
]
for cn,en,em,dangers,actions,img_h,hc,extra in weathers_data:
    # Brainstorm
    s=danger_safe_slide(cn,en,em,cn,en,dangers,actions,img_h,hc,brainstorm=True)
    n+=1;pn(s,n)
    notes(s,f"{em} {cn} brainstorm (2 分钟):\n• 不给答案 — 让学生自己想!\n• 老师把答案写白板上\n• 收集完 → 翻下一页一起看参考")
    # Reveal
    s=danger_safe_slide(cn,en,em,cn,en,dangers,actions,img_h,hc,brainstorm=False)
    n+=1;pn(s,n)
    notes(s,f"{em} {cn} 答案 (2 分钟):\n• 一起看标准答案\n• 对比学生 brainstorm 出来的\n• 表扬想到的, 补充没想到的\n• 提醒: {extra}")

# ===== TOPIC 2: ANIMALS =====
# 9. ANIMALS intro — image-driven hook
s=ns();bg(s,CREAM);hb(s,"🐾 野外的动物  Wild Animals",PINE)
ib(s,0.5,1.0,9.0,3.7,"📷 老师在这里插入野生动物图片 / Teacher: insert wildlife photo")
teacher_student_bar(s,4.85,
    "在野外, 你可能看到哪些动物?",
    "Think-Pair-Share: 我见过 / 想过 ___ 。")
n+=1;pn(s,n)
notes(s,"🐾 (3 分钟):\n• 这 3 条规则适用于所有野生动物\n• 让学生喊 3 次, 加深记忆\n• 接下来一个一个学具体动物应对")

# --- ANIMALS (brainstorm → reveal pair for each) ---
animals_data=[
    ("蛇","Snake","🐍",
        ["可能有毒","咬人很危险"],
        ["停下 — 不靠近","慢慢绕开","告诉大人"],
        "草丛里的蛇",DANGER,
        "不要戳 / 抓 — 蛇受惊才会咬"),
    ("熊","Bear","🐻",
        ["很大, 跑得比人快","可能攻击人"],
        ["不跑! 不尖叫!","慢慢后退","告诉大人"],
        "森林里的熊",BROWN,
        "反直觉: 不跑! 跑会让熊追! 慢慢后退 + 大声缓慢说话"),
    ("蜜蜂 / 黄蜂","Bee / Wasp","🐝",
        ["会蛰人","有毒, 会肿"],
        ["不挥手拍打!","慢慢离开","用衣服护住头"],
        "蜜蜂飞舞",CAUTION,
        "挥手反而让蜜蜂觉得受攻击 → 蛰你"),
    ("野狗","Wild Dog","🐕",
        ["可能咬人","成群更危险"],
        ["不跑 / 不尖叫","不盯着看","慢慢离开 + 叫大人"],
        "野狗",BROWN,
        "盯着看 = 挑衅; 转身跑 = 引追"),
    ("鹿 / 松鼠","Deer / Squirrel","🦌",
        ["看起来可爱, 但也是野生!","可能咬 / 抓人"],
        ["不摸 / 不喂","远远观察","用眼睛看, 不用手"],
        "森林里的小鹿和松鼠",SAFE,
        "投喂会改变动物习惯, 对它们不好"),
]
for cn,en,em,dangers,actions,img_h,hc,extra in animals_data:
    # Brainstorm
    s=danger_safe_slide(cn,en,em,cn,en,dangers,actions,img_h,hc,brainstorm=True)
    n+=1;pn(s,n)
    notes(s,f"{em} {cn} brainstorm (2 分钟):\n• 不给答案 — 让学生自己想!\n• 关键问题: 看到 {cn} — 你害怕吗? 你会怎么做?\n• 收集答案 → 翻下页看标准答案")
    # Reveal
    s=danger_safe_slide(cn,en,em,cn,en,dangers,actions,img_h,hc,brainstorm=False)
    n+=1;pn(s,n)
    notes(s,f"{em} {cn} 答案 (2 分钟):\n• 对比学生 brainstorm 答案\n• 表扬想到的\n• 重点: {extra}")

# ===== TOPIC 3: DON'T TOUCH/EAT — 10-slide expanded section =====
# A. TITLE PAGE — 野外安全
s=ns();bg(s,DANGER)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.0),W,Inches(2.5))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.55,8,0.7,"🌲 野外安全:",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.20,8,0.7,"这些不能碰!",sz=44,b=True,c=CAUTION,a=PP_ALIGN.CENTER)
tb(s,1,2.20,8,0.5,"Outdoor Safety: Don't Touch These!",sz=20,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,1,2.85,8,0.6,"⭐ 不认识的, 不吃 · 不碰 · 不靠近",sz=22,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,1,3.50,8,0.4,"If you don't know it — DON'T eat, touch, or go near it!",sz=13,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.80,8,0.4,"野外生存与探险 · Wilderness Safety",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「我们已经学了天气 + 动物 — 今天还要学一件事! 」\n• 「在野外 — 有很多东西看起来漂亮, 但是不能碰!」")

# B. WHY BE CAREFUL — predict question
s=ns();bg(s,CREAM);hb(s,"🤔 为什么要小心?  Why Be Careful?",DANGER)
tb(s,0.4,0.85,9.2,0.4,"在野外, 你觉得哪些东西不能碰?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"What things in the wild shouldn't we touch?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 1 big image placeholder showing collage of items
ib(s,0.4,1.65,9.2,2.85,"📷 老师插入图片 (野外物品集合: 蘑菇 / 野果 / 植物 / 带刺 / 垃圾) / Teacher: collage of outdoor items")
teacher_student_bar(s,4.65,
    "你觉得什么东西不能碰?",
    "Think-Pair-Share: 我觉得 ___ 不能碰。")
n+=1;pn(s,n)
notes(s,"hook (2-3 分钟):\n• 让学生想 — 不给答案!\n• 收集答案在白板上\n• 引出关键概念: 「看起来安全, 其实危险!」\n• 接下来一个一个看")

# C. 🍄 不要吃 — 毒蘑菇
s=ns();bg(s,CREAM);hb(s,"🍄 不要吃 — 毒蘑菇  Don't Eat: Mushrooms",DANGER)
ib(s,0.4,1.0,4.4,3.6,"📷 老师插入毒蘑菇真实图片 / Real photo: poisonous mushroom")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
tb(s,5.15,1.10,4.4,0.4,"⚠️ 危险  Why Dangerous",sz=14,b=True,c=DANGER)
tb(s,5.15,1.55,4.4,0.45,"·  看起来漂亮, 但可能有毒",sz=14,b=True,c=DARK)
tb(s,5.15,2.05,4.4,0.45,"·  吃了会生病",sz=14,b=True,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(2.85),Inches(4.6),Inches(1.75))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SAFE;sh2.line.width=Pt(2.5)
tb(s,5.15,2.95,4.4,0.4,"✅ 规则  Rule",sz=14,b=True,c=SAFE)
tb(s,5.15,3.40,4.4,0.6,"不认识的蘑菇 — 不吃!",sz=22,b=True,c=DARK)
tb(s,5.15,4.05,4.4,0.4,"Don't eat unknown mushrooms!",sz=11,c=GRAY)
teacher_student_bar(s,4.75,
    "蘑菇漂亮 — 可以吃吗?",
    "我不吃! 不认识的蘑菇不吃!")
n+=1;pn(s,n)
notes(s,"🍄 毒蘑菇 (2 分钟):\n• 强调反直觉: 漂亮 ≠ 安全\n• 真正的毒蘑菇: 鹅膏 / 红伞伞 / 白杆杆\n• 美国常见毒蘑菇: Death Cap / Destroying Angel\n• 「就算看起来像你在超市见过 — 不一定一样!」")

# D. 🫐 不要吃 — 野果
s=ns();bg(s,CREAM);hb(s,"🫐 不要吃 — 野果  Don't Eat: Wild Berries",DANGER)
ib(s,0.4,1.0,4.4,3.6,"📷 老师插入野果真实图片 (像蓝莓但有毒) / Wild berries that look like blueberries")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
tb(s,5.15,1.10,4.4,0.4,"⚠️ 危险  Why Dangerous",sz=14,b=True,c=DANGER)
tb(s,5.15,1.55,4.4,0.45,"·  看起来像蓝莓",sz=14,b=True,c=DARK)
tb(s,5.15,2.05,4.4,0.45,"·  但可能有毒, 会中毒",sz=14,b=True,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(2.85),Inches(4.6),Inches(1.75))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SAFE;sh2.line.width=Pt(2.5)
tb(s,5.15,2.95,4.4,0.4,"✅ 规则  Rule",sz=14,b=True,c=SAFE)
tb(s,5.15,3.40,4.4,0.6,"不是家里来的 — 不吃!",sz=22,b=True,c=DARK)
tb(s,5.15,4.05,4.4,0.4,"Not from home? Don't eat!",sz=11,c=GRAY)
teacher_student_bar(s,4.75,
    "看起来像蓝莓 — 可以吃吗?",
    "我不吃! 不是家里的不吃!")
n+=1;pn(s,n)
notes(s,"🫐 野果 (2 分钟):\n• 「看起来像蓝莓」 — 但可能是颠茄 / 龙葵 (有毒!)\n• 美国警告: pokeberry / nightshade — 看起来美味但有毒\n• 即使鸟吃 → 人不一定能吃")

# E. 🌱 不要碰 — 陌生植物
s=ns();bg(s,CREAM);hb(s,"🌱 不要碰 — 陌生植物  Don't Touch: Strange Plants",DANGER)
ib(s,0.4,1.0,4.4,3.6,"📷 老师插入植物图片 (毒漆树 / poison ivy) / Real photo: poison ivy")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
tb(s,5.15,1.10,4.4,0.4,"⚠️ 危险  Why Dangerous",sz=14,b=True,c=DANGER)
tb(s,5.15,1.55,4.4,0.45,"·  皮肤会痒, 起红疹",sz=14,b=True,c=DARK)
tb(s,5.15,2.05,4.4,0.45,"·  例: 毒漆树 (Poison Ivy)",sz=14,b=True,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(2.85),Inches(4.6),Inches(1.75))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SAFE;sh2.line.width=Pt(2.5)
tb(s,5.15,2.95,4.4,0.4,"✅ 规则  Rule",sz=14,b=True,c=SAFE)
tb(s,5.15,3.40,4.4,0.6,"不认识的植物 — 不碰!",sz=22,b=True,c=DARK)
tb(s,5.15,4.05,4.4,0.4,"Don't touch unknown plants!",sz=11,c=GRAY)
teacher_student_bar(s,4.75,
    "看到不认识的叶子 — 怎么办?",
    "我不碰! 走开!")
n+=1;pn(s,n)
notes(s,"🌱 陌生植物 (2 分钟):\n• 美国常见: Poison Ivy / Poison Oak / Poison Sumac\n• 口诀: 「Leaves of three — let it be!」 (3 片叶子 — 别碰!)\n• 接触后皮肤会起红疹, 痒一周\n• 万一碰了: 用大量清水 + 肥皂洗")

# F. 🌵 不要碰 — 带刺植物
s=ns();bg(s,CREAM);hb(s,"🌵 不要碰 — 带刺植物  Don't Touch: Thorny Plants",DANGER)
ib(s,0.4,1.0,4.4,3.6,"📷 老师插入带刺植物图片 (仙人掌 / 玫瑰丛) / Real photo: cactus / thorns")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
tb(s,5.15,1.10,4.4,0.4,"⚠️ 危险  Why Dangerous",sz=14,b=True,c=DANGER)
tb(s,5.15,1.55,4.4,0.45,"·  扎手, 很疼",sz=14,b=True,c=DARK)
tb(s,5.15,2.05,4.4,0.45,"·  可能流血 / 感染",sz=14,b=True,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(2.85),Inches(4.6),Inches(1.75))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SAFE;sh2.line.width=Pt(2.5)
tb(s,5.15,2.95,4.4,0.4,"✅ 规则  Rule",sz=14,b=True,c=SAFE)
tb(s,5.15,3.40,4.4,0.6,"看到刺 — 离远一点!",sz=22,b=True,c=DARK)
tb(s,5.15,4.05,4.4,0.4,"See thorns? Stay away!",sz=11,c=GRAY)
teacher_student_bar(s,4.75,
    "看到刺 — 怎么办?",
    "我走开! 不碰!")
n+=1;pn(s,n)
notes(s,"🌵 带刺植物 (2 分钟):\n• 仙人掌 / 玫瑰 / 荨麻 / 蓝蓟\n• 万一扎到: 不用手拔 (会留一节)\n• 用镊子 / 胶带慢慢拔出\n• 然后洗 + 消毒")

# G. 🍼 不要碰 — 不明物品 / 垃圾
s=ns();bg(s,CREAM);hb(s,"🍼 不要碰 — 不明物品 / 垃圾  Don't Touch: Strange Items",DANGER)
ib(s,0.4,1.0,4.4,3.6,"📷 老师插入图片 (野外垃圾 / 旧瓶子 / 旧绳子) / Trash, old bottles in nature")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
tb(s,5.15,1.10,4.4,0.4,"⚠️ 危险  Why Dangerous",sz=14,b=True,c=DANGER)
tb(s,5.15,1.55,4.4,0.45,"·  里面可能有危险物品",sz=14,b=True,c=DARK)
tb(s,5.15,2.05,4.4,0.45,"·  可能有化学物质",sz=14,b=True,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(2.85),Inches(4.6),Inches(1.75))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SAFE;sh2.line.width=Pt(2.5)
tb(s,5.15,2.95,4.4,0.4,"✅ 规则  Rule",sz=14,b=True,c=SAFE)
tb(s,5.15,3.40,4.4,0.6,"不是你的 — 不碰!",sz=22,b=True,c=DARK)
tb(s,5.15,4.05,4.4,0.4,"Not yours? Don't touch!",sz=11,c=GRAY)
teacher_student_bar(s,4.75,
    "看到一个旧瓶子 / 注射器 / 玻璃 — 怎么办?",
    "我不碰! 告诉大人!")
n+=1;pn(s,n)
notes(s,"🍼 不明物品 (2 分钟):\n• 包括: 旧瓶子 / 注射器 / 玻璃片 / 锈钉 / 旧绳子\n• 化学物质可能烧伤皮肤\n• 玻璃 / 注射器可能割伤 / 感染\n• 永远告诉大人, 让大人处理")

# H. 🧠 INTERACTIVE QUIZ
s=ns();bg(s,CREAM);hb(s,"🧠 小测试  Quick Quiz!",CALM)
tb(s,0.4,0.85,9.2,0.4,"举手 / 大声喊! YES 还是 NO?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Raise hands or shout: YES or NO?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Quiz 1: 哪个可以吃?
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.5),Inches(2.85))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=CALM;sh.line.width=Pt(2.5)
tb(s,0.55,1.75,4.3,0.4,"Q1. 哪个可以吃?",sz=15,b=True,c=CALM)
tb(s,0.55,2.15,4.3,0.3,"Which one is safe to eat?",sz=10,c=GRAY)
# Two options
opt1a=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(2.55),Inches(2.0),Inches(1.75))
opt1a.fill.solid();opt1a.fill.fore_color.rgb=WHITE;opt1a.line.color.rgb=DANGER;opt1a.line.width=Pt(2)
tb(s,0.55,2.65,2.0,1.0,"🍄",sz=72,a=PP_ALIGN.CENTER)
tb(s,0.55,3.75,2.0,0.4,"野蘑菇",sz=14,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,0.55,4.10,2.0,0.3,"Wild mushroom",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
opt1b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.75),Inches(2.55),Inches(2.0),Inches(1.75))
opt1b.fill.solid();opt1b.fill.fore_color.rgb=WHITE;opt1b.line.color.rgb=SAFE;opt1b.line.width=Pt(2)
tb(s,2.75,2.65,2.0,1.0,"🍎",sz=72,a=PP_ALIGN.CENTER)
tb(s,2.75,3.75,2.0,0.4,"家里苹果",sz=14,b=True,c=SAFE,a=PP_ALIGN.CENTER)
tb(s,2.75,4.10,2.0,0.3,"Apple from home",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Quiz 2: 哪个可以碰?
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.65),Inches(4.5),Inches(2.85))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=CALM;sh2.line.width=Pt(2.5)
tb(s,5.25,1.75,4.3,0.4,"Q2. 哪个可以碰?",sz=15,b=True,c=CALM)
tb(s,5.25,2.15,4.3,0.3,"Which one is safe to touch?",sz=10,c=GRAY)
opt2a=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.25),Inches(2.55),Inches(2.0),Inches(1.75))
opt2a.fill.solid();opt2a.fill.fore_color.rgb=WHITE;opt2a.line.color.rgb=DANGER;opt2a.line.width=Pt(2)
tb(s,5.25,2.65,2.0,1.0,"🌵",sz=72,a=PP_ALIGN.CENTER)
tb(s,5.25,3.75,2.0,0.4,"仙人掌",sz=14,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,5.25,4.10,2.0,0.3,"Cactus",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
opt2b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.45),Inches(2.55),Inches(2.0),Inches(1.75))
opt2b.fill.solid();opt2b.fill.fore_color.rgb=WHITE;opt2b.line.color.rgb=SAFE;opt2b.line.width=Pt(2)
tb(s,7.45,2.65,2.0,1.0,"🧸",sz=72,a=PP_ALIGN.CENTER)
tb(s,7.45,3.75,2.0,0.4,"小熊玩偶",sz=14,b=True,c=SAFE,a=PP_ALIGN.CENTER)
tb(s,7.45,4.10,2.0,0.3,"Teddy bear",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.65,
    "举手! Q1: 蘑菇还是苹果? Q2: 仙人掌还是玩偶?",
    "✅ 苹果! ✅ 玩偶!")
n+=1;pn(s,n)
notes(s,"小测试 (3-4 分钟):\n• 让学生大声喊答案\n• Q1 答案: 苹果 (家里来的才安全)\n• Q2 答案: 玩偶 (没刺, 没毒)\n• 表扬答对的\n• 也可以让学生自己出类似题目 (野苺 vs 草莓; 不明瓶子 vs 水壶)")

# I. 三个不 — RULES SUMMARY
s=ns();bg(s,CREAM);hb(s,"⭐ 三个 「不」 规则  3 \"Don't\" Rules",SAFE)
tb(s,0.4,0.85,9.2,0.4,"记住这 3 个 — 你就安全!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Remember these 3 — and you're safe!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
rules=[("🚫","不吃","Don't eat",DANGER),
       ("🚫","不碰","Don't touch",DANGER),
       ("🚫","不靠近","Don't go near",DANGER)]
for i,(em,cn,en,cl) in enumerate(rules):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(3.0),Inches(2.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.75,2.8,1.2,em,sz=80,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.05,2.8,0.6,cn,sz=28,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.70,2.8,0.4,en,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.10,2.8,0.4,"❌",sz=22,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.65,
    "三个不是什么?",
    "全班大声喊: 不吃! 不碰! 不靠近!")
n+=1;pn(s,n)

# J. 一句话记住 — KEY TAKEAWAY
s=ns();bg(s,CAUTION)
tb(s,1,0.9,8,0.6,"🎯 一句话记住!",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.55,8,0.4,"Remember in One Sentence!",sz=14,c=WHITE,a=PP_ALIGN.CENTER)
# Big quote bubble
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(2.30),Inches(9.0),Inches(2.30))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(4)
tb(s,0.5,2.60,9.0,0.7,"「不认识的,」",sz=32,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,0.5,3.30,9.0,0.7,"不吃 · 不碰 · 不靠近!」",sz=32,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,0.5,4.05,9.0,0.4,"\"If you don't know it — DON'T eat, touch, or go near it!\"",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,4.85,8,0.4,"全班一起念 3 次!  Class — say it 3 times!",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"总结 (1-2 分钟):\n• 一起大声念 3 次\n• 越念越大声\n• 最后一次加动作: 摇头 + 摆手\n• 「这一句 — 记住了, 你就是安全探险家!」")
notes(s,"答案揭晓 (3 分钟):\n• 一个一个指着说\n• 对比学生 brainstorm 出来的\n• 不要详细讲哪个有毒 — 太多了\n• 重点: 陌生 = 不碰\n• 「就算看起来像你在家里见过的 — 不一定一样!」")

# ===== TOPIC 4: LOST =====
# 17. LOST intro — predict (think first)
s=ns();bg(s,CREAM);hb(s,"迷路了怎么办?  What if You Get Lost?",CALM)
tb(s,0.4,0.78,9.2,0.30,"想一想 — 你会怎么做? Think — what would YOU do?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# Big picture placeholder on LEFT — lost kid scene
ib(s,0.4,1.15,4.4,3.65,"老师插入图片: 迷路的小朋友 (Lost child scene)")
# 4 brainstorm boxes on RIGHT — 2x2 grid, no question marks
for i in range(4):
    col=i%2;row=i//2
    x=5.0+col*2.35;y=1.15+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(1.70))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CALM;sh.line.width=Pt(2.5)
    pill(s,x+0.12,y+0.12,0.45,0.40,str(i+1),CALM,sz=16)
    tb(s,x+0.05,y+0.70,2.10,0.40,f"想法 {i+1}",sz=15,b=True,c=CALM,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.10,2.10,0.30,f"Idea {i+1}",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.95,
    "你一个人在森林里, 找不到大人 — 怎么办?",
    "Think-Pair-Share — 我会 ___ 。",color=CALM)
n+=1;pn(s,n)
notes(s,"hook (3 分钟):\n• 不要给答案 — 让学生想\n• 收集所有想法在白板上, 写到 4 个 想法 框里\n• 常见错误想法: 「跑回去找」「乱跑找路」「自己穿过森林」 — 都不对!\n• 引出 6 步安全做法\n• 老师可以提前在图片框插入: 迷路小朋友独自在森林路口 / 找不到大人 等场景照片")

# 18. LOST — Right or Wrong? (activate prior knowledge — quiz first)
s=ns();bg(s,CREAM);hb(s,"对还是错?  Right or Wrong — Lost?",CALM)
tb(s,0.4,0.78,9.2,0.30,"迷路了, 这样做对吗? — Is this what you should do?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
quiz=[
    ("1","我跑回去找路。","I run back to find the path."),
    ("2","我留在原地。","I stay where I am."),
    ("3","我下山谷找路。","I go down the valley."),
    ("4","我大声喊!","I shout out loud!"),
]
for i,(num,cn,en) in enumerate(quiz):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.18+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CALM;sh.line.width=Pt(2.5)
    pill(s,x+0.12,y+0.12,0.45,0.40,num,CALM,sz=16)
    ib(s,x+0.15,y+0.62,1.50,1.05,"老师插入图片")
    tb(s,x+1.80,y+0.55,2.65,0.55,cn,sz=20,b=True,c=DARK)
    tb(s,x+1.80,y+1.15,2.65,0.30,en,sz=10,c=GRAY)
teacher_student_bar(s,4.95,
    "对? 还是错? 全班一起举手!",
    "对 → 拍手一次; 错 → 摇摇头",color=CALM)
n+=1;pn(s,n)
notes(s,"对错快答 (4-5 分钟):\n• 不要立刻给答案 — 让学生先猜!\n• 一道一道, 全班举手投票\n• 答案:\n  - 1. 我跑回去找路。 → 错! 越跑离大人越远\n  - 2. 我留在原地。 → 对! 留原地大人才找得到\n  - 3. 我下山谷找路。 → 错! 山谷地形危险, 水声盖过呼救\n  - 4. 我大声喊! → 对! 也可以吹哨子\n• 引导讨论: 「为什么是错的?」「为什么是对的?」\n• 学生猜完 → 翻下一页看正确做法\n• 老师可以提前在每个图片框插入对应的真实场景照片")

# 19. LOST — What to Do When Lost (visual answer — 4 essentials)
s=ns();bg(s,CREAM);hb(s,"迷路了怎么办?  What to Do When Lost",CALM)
tb(s,0.4,0.78,9.2,0.30,"记住这 4 件事 — Remember these 4 things",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
actions=[
    ("停","Stop","站住, 不要跑","Don't run",DANGER),
    ("留","Stay","留在原地","Stay in place",CAUTION),
    ("喊","Call","大声喊 / 吹哨子","Shout / whistle",SUN),
    ("等","Wait","等大人来","Wait for adults",SAFE),
]
for i,(big,en_top,cn_desc,en_desc,cl) in enumerate(actions):
    x=0.45+i*2.30;y=1.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(3.70))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    ib(s,x+0.10,y+0.15,2.00,1.70,"老师插入图片")
    tb(s,x+0.05,y+1.95,2.10,1.00,big,sz=60,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+3.05,2.10,0.30,cn_desc,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+3.40,2.10,0.25,en_desc,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.95,
    "一起做 — 停! 留! 喊! 等!",
    "全班一起喊 3 次, 加动作!",color=CALM)
n+=1;pn(s,n)
notes(s,"正确做法 (3-4 分钟):\n• 揭晓答案 — 对应上一页的对错题\n• 4 件事 (越简单越记得住):\n  - 停 — 看到自己迷路, 立刻停下 (深呼吸 3 次)\n  - 留 — 留在原地, 不乱跑 (大人才找得到)\n  - 喊 — 大声呼叫 / 吹哨子 (3 长信号, 等 30 秒)\n  - 等 — 等大人来, 不要自己穿森林\n• 关键反直觉点: 不要跑! 不要下山谷!\n• 全班一起喊 3 次: 「停 — 留 — 喊 — 等!」加动作\n• 联系 Day 2 装包: 哨子是必备\n• 老师可以提前准备 4 张真实照片插入图片框 (停下 / 留原地 / 吹哨子 / 跟大人走)")

# 19b. LOST — S.T.O.P. principle (deeper dive)
s=ns();bg(s,CREAM);hb(s,"🛑 S.T.O.P. 法则  The S.T.O.P. Rule",CALM)
ib(s,0.4,0.95,4.4,3.7,"📷 老师插入图片: 小朋友停下来思考 / Child stopping to think")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(0.95),Inches(4.6),Inches(3.7))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=CALM;sh.line.width=Pt(2.5)
tb(s,5.15,1.05,4.4,0.4,"🌍 国际通用法则",sz=13,b=True,c=CALM)
stops=[("S","Stop","停下来",DANGER),
       ("T","Think","想一想",CAUTION),
       ("O","Observe","看一看",NAVY),
       ("P","Plan","想办法",SAFE)]
for i,(ltr,en,cn,cl) in enumerate(stops):
    y=1.50+i*0.78
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.15),Inches(y),Inches(4.30),Inches(0.70))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=cl;sh2.line.width=Pt(1.5)
    pill(s,5.25,y+0.13,0.55,0.45,ltr,cl,sz=18)
    tb(s,5.90,y+0.06,1.4,0.30,en,sz=11,c=GRAY)
    tb(s,5.90,y+0.30,3.4,0.40,cn,sz=18,b=True,c=cl)
teacher_student_bar(s,4.80,
    "你一个人在森林 — S.T.O.P. 是什么?",
    "全班喊: 停! 想! 看! 想办法!")
n+=1;pn(s,n)
notes(s,"S.T.O.P. 法则 (3-4 分钟):\n• 这是国际野外求生通用原则\n• 不要慌, 不要跑\n• S - Stop: 停下来 (深呼吸 3 次)\n• T - Think: 想一想上次看到的路 / 大人\n• O - Observe: 看一看周围 — 有什么标志? 听到什么?\n• P - Plan: 想办法 — 留原地? 找掩体? 用哨子?\n• 让学生自己演练一遍")

# 19c. LOST — How to signal for help (求救信号)
s=ns();bg(s,CREAM);hb(s,"怎么求救?  How to Call for Help",CALM)
tb(s,0.4,0.78,9.2,0.30,"4 种求救信号  4 Ways to Signal",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
signals=[
    ("打电话","Call",["110 / 119","911 (美国)"],CALM),
    ("吹哨子","Whistle",["3 长哨!","等 30 秒","再 3 长"],SUN),
    ("鲜艳衣服","Bright clothes",["放空旷处","让救援看见"],DANGER),
    ("「3」是信号","Rule of 3",["3 堆石头","3 团烟","3 个任何"],CAUTION),
]
for i,(cn,en,details,cl) in enumerate(signals):
    x=0.4+i*2.35;y=1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(3.35))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    # Large image placeholder (top — for real photo)
    ib(s,x+0.10,y+0.10,2.00,1.50,"老师插入图片")
    # Chinese title
    tb(s,x+0.05,y+1.70,2.10,0.35,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    # English title
    tb(s,x+0.05,y+2.05,2.10,0.25,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # Details
    yy=y+2.35
    for d in details:
        tb(s,x+0.10,yy,2.0,0.26,f"·  {d}",sz=11,c=DARK);yy+=0.25
teacher_student_bar(s,4.60,
    "记住! 「3」 = 国际求救信号",
    "3 长哨! 3 堆石头! 3 团烟!",color=CALM)
n+=1;pn(s,n)
notes(s,"求救信号 (3-4 分钟):\n• 国际通用: 「3」 是 SOS\n• 哨子: 3 长哨, 等 30 秒, 再 3 长\n• 视觉: 鲜艳衣服 (红/黄/橙) 放空旷处\n• 烟雾: 3 堆火 / 3 团烟 (老师监督)\n• 电话: 110 (中国) / 119 (中国/日本) / 911 (美国/加拿大)\n• 让学生拿哨子练一次 (Day 4 Project 3 工坊)")

# 19d. LOST — Find direction (without compass)
s=ns();bg(s,CREAM);hb(s,"🧭 怎么辨别方向?  Find Your Direction",NAVY)
ib(s,0.4,0.95,4.4,3.7,"📷 老师插入图片: 太阳 / 树木 / 指南针 / Sun / tree / compass")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(0.95),Inches(4.6),Inches(3.7))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=NAVY;sh.line.width=Pt(2.5)
tb(s,5.15,1.05,4.4,0.4,"3 个方法  3 Methods",sz=14,b=True,c=NAVY)
methods=[
    ("🧭","指南针","红针指北 (Red = N)",DANGER),
    ("☀️","太阳 + 手表","时针对太阳, 时针与 12 的中间 = 南 (北半球)",SUN),
    ("🌳","树 + 青苔","南面叶子茂; 北面长青苔",PINE),
]
for i,(em,cn,detail,cl) in enumerate(methods):
    y=1.55+i*1.00
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.15),Inches(y),Inches(4.30),Inches(0.90))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=cl;sh2.line.width=Pt(1.5)
    tb(s,5.25,y+0.18,0.6,0.55,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,5.90,y+0.05,3.4,0.30,cn,sz=15,b=True,c=cl)
    tb(s,5.90,y+0.35,3.4,0.50,detail,sz=10,c=DARK)
teacher_student_bar(s,4.80,
    "没有指南针 — 怎么找方向?",
    "看太阳 / 看树 — 但最准的是指南针!")
n+=1;pn(s,n)
notes(s,"辨别方向 (3-4 分钟):\n• 联系 Day 3 学过的内容\n• 1. 指南针: 最准! 红针 = 北\n• 2. 太阳 + 手表 (G2+ 才教):\n  - 北半球: 把时针对太阳\n  - 时针与 12 点之间的角 — 平分线 = 南\n• 3. 树木 / 青苔: 南面阳光多 → 茂盛\n  - 北面阴 → 长青苔 (但不100% 准)\n• 强调: 自然线索不一定准 — 最好用指南针")

# 19e. LOST — What NOT to do
s=ns();bg(s,CREAM);hb(s,"❌ 千万不要!  NEVER Do These!",DANGER)
tb(s,0.4,0.85,9.2,0.4,"两个致命错误 — 不要做!",sz=18,b=True,c=DANGER,a=PP_ALIGN.CENTER)
dont=[
    ("🏃‍♂️","不要慌跑","Don't run blindly","只会浪费体力, 越跑越远离大人",
     "孩子在森林里乱跑"),
    ("⛰️","不要下山谷 / 溪谷","Don't go down valleys",
     "山谷陡峭危险, 水声会盖过你的喊声",
     "陡峭山谷与溪流"),
]
for i,(em,cn,en,why,img_h) in enumerate(dont):
    x=0.4+i*4.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.40),Inches(4.5),Inches(3.30))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(3)
    ib(s,x+0.15,1.55,4.20,1.65,f"📷 {img_h}")
    tb(s,x+0.15,3.30,4.20,0.5,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,3.85,4.20,0.45,cn,sz=18,b=True,c=DANGER,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,4.30,4.20,0.30,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    # Why text moved to a separate small caption underneath card
teacher_student_bar(s,4.85,
    "迷路了 — 你会跑吗?",
    "不跑! 不下山谷! 留原地, 等大人!")
n+=1;pn(s,n)
notes(s,"反直觉重点 (3 分钟):\n• ❌ 不要慌乱跑:\n  - 体力浪费\n  - 越跑越偏离原路\n  - 大人找你也难\n• ❌ 不要下山谷 / 溪谷:\n  - 山谷地形陡峭, 危险\n  - 水流声音会盖过你的喊声 (大人听不到)\n  - 万一受伤更难救\n• 这两点都是反直觉 — 强调!\n• 联系 S.T.O.P. 的第一步: STOP")

# 19. LOST — sentence frame practice
s=ns();bg(s,CREAM);hb(s,"💬 句型练习  Sentence Practice",CALM)
tb(s,0.4,0.85,9.2,0.4,"我看到 ___ , 我不 ___ , 我会 ___ 。",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.3,"I see ___ , I don't ___ , I will ___ .",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
examples=[
    ("🐍","蛇","跑","慢慢绕开"),
    ("🐻","熊","跑","慢慢后退"),
    ("🍄","毒蘑菇","吃","告诉大人"),
    ("❓","迷路了","乱跑","留在原地, 大声喊"),
]
for i,(em,what,wrong,right) in enumerate(examples):
    y=1.75+i*0.78
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.70))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CALM;sh.line.width=Pt(2)
    tb(s,0.55,y+0.13,0.7,0.4,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.20,2.4,0.4,f"我看到 {what}",sz=14,b=True,c=DARK)
    tb(s,3.85,y+0.20,2.4,0.4,f"我不 {wrong}",sz=14,b=True,c=DANGER)
    tb(s,6.50,y+0.20,3.0,0.4,f"我会 {right}",sz=14,b=True,c=SAFE)
n+=1;pn(s,n)
notes(s,"句型练习 (3 分钟):\n• 让学生一行一行读\n• 然后各自编一句\n• 这个句型是 Session 3 表演的核心")

# ===== SESSION 2 =====
# 20. SESSION 2 DIVIDER
s=div("Session 2  下午","🔄 复习 + 我会认 4 / 我会写 2",CAUTION,"📖");n+=1;pn(s,n)

# 21. REVIEW — 4 topics quick recap
s=ns();bg(s,CREAM);hb(s,"🔄 上午学了什么?  Review",CAUTION)
tb(s,0.4,0.85,9.2,0.4,"我们学了 4 种可能遇到的危险 — 你还记得吗?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
topics=[("🌤️","天气","Weather","晴/雨/风/冷",CAUTION),
        ("🐾","动物","Animals","蛇/熊/蜂/狗",PINE),
        ("🚫","不碰","Don't touch","蘑菇/野果/陌生物",DANGER),
        ("❓","迷路","Lost","6 步: 停/留/喊/哨/标/等",CALM)]
for i,(em,cn,en,detail,cl) in enumerate(topics):
    col=i%2;row=i//2
    x=0.4+col*4.65;y=1.45+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(1.50))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.10,1.0,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.10,3.2,0.4,cn,sz=22,b=True,c=cl)
    tb(s,x+1.20,y+0.55,3.2,0.3,en,sz=11,c=GRAY)
    tb(s,x+1.20,y+0.95,3.2,0.4,detail,sz=12,c=DARK)
teacher_student_bar(s,4.85,
    "哪一个你最印象深刻?",
    "我记得 ___, 因为 ___ 。")
n+=1;pn(s,n)

# 22-25. WORD CARDS — 我会认 (4 words)
read_data=[
    ("天气","tiān qì","Weather","☁️","今天的天气很好。"),
    ("危险","wēi xiǎn","Danger","⚠️","看到蛇 — 危险! 慢慢绕开。"),
    ("迷路","mí lù","Lost","🗺️","迷路了 — 不跑, 留在原地!"),
    ("食物","shí wù","Food","🍱","陌生的食物不能吃。"),
]
for w,py,en,em,sent in read_data:
    s=word_card_read(w,py,en,em,sent,color=DANGER);n+=1;pn(s,n)

# 26-27. WORD CARDS — 我会写 (2 — 天气 / 危险)
write_data=[
    ("天气","tiān qì","Weather", 8,"8 笔: 天 (一+大) + 气 (气字旁 4 笔) · Sky + breath"),
    ("危险","wēi xiǎn","Danger", 13,"13 笔: 危 (㔾上面) + 险 (阝+佥) · Cliff + steep"),
]
for w,py,en,sc,hint in write_data:
    s=word_card_write(w,py,en,sc,hint,color=DANGER);n+=1;pn(s,n)

# 28. REVIEW GAME — Baamboozle placeholder
s=ns();bg(s,CREAM);hb(s,"🎮 复习游戏  Review · Baamboozle",CAUTION)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(1.30),Inches(7.0),Inches(3.3))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=CAUTION;sh.line.width=Pt(2.5)
tb(s,1.5,1.85,7.0,0.6,"🎲",sz=80,a=PP_ALIGN.CENTER)
tb(s,1.5,2.85,7.0,0.5,"老师在这里放 Baamboozle 链接",sz=20,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,1.5,3.40,7.0,0.4,"Teacher: paste your Baamboozle link here",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1.5,3.90,7.0,0.4,"🔗 https://www.baamboozle.com/...",sz=14,b=True,c=CAUTION,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.80,
    "我们来玩复习游戏!",
    "分组抢答 — 答对得分!")
n+=1;pn(s,n)
notes(s,"复习游戏 (10-15 分钟):\n• 题目涵盖: 4 种天气 / 5 种动物 / 3 个不能碰规则 / 迷路6步\n• 老师提前在 Baamboozle 创建\n• 全班分 2-3 组")

# ===== SESSION 3 =====
# 29. SESSION 3 DIVIDER
s=div("Session 3  下午","🛠️ 表演 + Poster + 更多活动",SAFE,"🎭");n+=1;pn(s,n)

# 30. PROJECTS OVERVIEW
s=ns();bg(s,CREAM);hb(s,"🛠️ 动手时间  Hands-On Time — 3 Activities",SAFE)
projects=[
    ("ACTIVITY 1","🎭 安全情景小剧场","Safety Skits","分组表演 4 个情景\nGroup performance",WARM,DANGER,"主项目 ⭐"),
    ("ACTIVITY 2","📓 安全小书 / Poster","Safety Booklet","4 页: 天气/动物/迷路/食物\n4 pages",RGBColor(0xFF,0xE0,0xB2),NAVY,"个人作品"),
    ("ACTIVITY 3","🆘 紧急哨子工坊","Whistle Workshop","DIY 应急哨子 + 训练\nDIY whistle + drill",RGBColor(0xDC,0xED,0xC8),SAFE,"加强 / Bonus"),
]
for i,(lbl,nm,en,d,bgc,cl,lvl) in enumerate(projects):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(3.1),Inches(4.15))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,2.9,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.4,2.9,0.6,nm,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.0,2.9,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.55,2.45,2.0,0.35,lvl,cl,sz=10)
    ib(s,x+0.2,2.95,2.8,1.2,"📷 示范")
    ls=d.split('\n')
    tf=tb(s,x+0.15,4.20,2.85,0.4,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for ln in ls[1:]:ap(tf,ln,sz=12,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# 31. ACTIVITY 1 — Safety Skits (intro: scenarios + roles)
s=ns();bg(s,CREAM);hb(s,"🎭 ACTIVITY 1: 安全情景小剧场  Safety Skits",DANGER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.40))
sh.fill.solid();sh.fill.fore_color.rgb=DANGER;sh.line.fill.background()
tb(s,0.4,0.99,4.2,0.35,"🎯 4 个情景 (抽签)  4 Scenarios",sz=13,b=True,c=WHITE)
scenes=[("🌧️","下雨了怎么办?","Rainy day"),
        ("🐍","看到蛇怎么办?","Saw a snake"),
        ("❓","迷路了怎么办?","Got lost"),
        ("🐝","被蜜蜂追怎么办?","Chased by bees")]
for i,(em,cn,en) in enumerate(scenes):
    y=1.45+i*0.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(4.2),Inches(0.50))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DANGER;sh.line.width=Pt(1.5)
    tb(s,0.5,y+0.10,0.6,0.35,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.15,y+0.05,3.0,0.3,cn,sz=14,b=True,c=DARK)
    tb(s,1.15,y+0.30,3.0,0.25,en,sz=10,c=GRAY)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(0.40))
sh2.fill.solid();sh2.fill.fore_color.rgb=SAFE;sh2.line.fill.background()
tb(s,4.95,0.99,4.65,0.35,"👥 每组 5 人 — 5 个角色",sz=13,b=True,c=WHITE)
roles=[("🧒","主角","Hero (问题主角)"),
       ("⚠️","危险","Danger (动物/天气)"),
       ("❌","错误做法","Wrong way"),
       ("✅","正确做法","Right way"),
       ("🎤","旁白","Narrator")]
for i,(em,cn,en) in enumerate(roles):
    y=1.45+i*0.42
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(y),Inches(4.65),Inches(0.38))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SAFE;sh.line.width=Pt(1.5)
    tb(s,5.05,y+0.05,0.5,0.30,em,sz=18,a=PP_ALIGN.CENTER)
    tb(s,5.55,y+0.04,1.6,0.30,cn,sz=13,b=True,c=SAFE)
    tb(s,7.20,y+0.04,2.4,0.30,en,sz=10,c=GRAY)
teacher_student_bar(s,4.85,
    "选角色 + 排练 5-10 分钟!",
    "想一想: 发生什么 → 错做法 → 对做法")
n+=1;pn(s,n)
notes(s,"安全表演 (15-20 分钟):\n• 分组: 5 人一组 (4 组共 20 人)\n• 抽 1 个情景\n• 每人选 1 个角色\n• 5-10 分钟讨论 + 排练\n• 道具极简 — 用动作 / 椅子 / 衣服即可")

# 32. ACTIVITY 1 — Performance flow + sentence frame
s=ns();bg(s,CREAM);hb(s,"🎬 表演流程 + 必说句子  Performance",DANGER)
flow=[("1️⃣","发生什么","What happened",CAUTION),
      ("2️⃣","错误做法","Wrong way",DANGER),
      ("3️⃣","正确做法","Right way!",SAFE)]
for i,(num,cn,en,cl) in enumerate(flow):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(3.0),Inches(1.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,2.8,0.5,num,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.55,2.8,0.4,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.00,2.8,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
for ax in [3.50,6.65]:
    tb(s,ax,1.55,0.30,0.5,"→",sz=24,b=True,c=GRAY,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(2.65),Inches(9.2),Inches(1.85))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=DANGER;sh.line.width=Pt(2.5)
tb(s,0.55,2.78,9.0,0.4,"💬 至少说一句!  At Least Say This:",sz=14,b=True,c=DANGER)
tb(s,0.55,3.30,9.0,0.7,"我看到了 ___,",sz=22,b=True,c=DARK)
tb(s,0.55,3.80,9.0,0.7,"我不 ___,  我会 ___ 。",sz=22,b=True,c=DARK)
tb(s,0.55,4.35,9.0,0.3,"I saw ___, I don't ___, I will ___.",sz=11,c=GRAY)
teacher_student_bar(s,4.65,
    "每组上台表演!",
    "其他组鼓掌 + 当评委!")
n+=1;pn(s,n)
notes(s,"表演展示 (15 分钟):\n• 每组 2-3 分钟表演\n• 表演完全班鼓掌\n• 老师简短评论: 「他们演得好在哪里?」\n• 收集共同学到的安全做法")

# 33. ACTIVITY 2 — Safety Booklet
s=ns();bg(s,CREAM);hb(s,"📓 ACTIVITY 2: 安全小书 / Poster  Safety Booklet",NAVY)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.40))
sh.fill.solid();sh.fill.fore_color.rgb=NAVY;sh.line.fill.background()
tb(s,0.4,0.99,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.4,"📄 小册子纸 / 打印模板",sz=13,c=DARK)
ap(tf,"   Booklet paper / printable template",sz=10,c=GRAY)
ap(tf,"🖍️ 彩笔 / 蜡笔  Markers / crayons",sz=13,c=DARK)
ap(tf,"✂️ 剪刀 + 胶水  Scissors + glue",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(0.40))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,4.95,0.99,4.65,0.35,"📑 4 页主题  4 Pages",sz=14,b=True,c=WHITE)
pages=[("☁️","天气 — 下雨怎么办?",CAUTION),
       ("🐍","危险 — 看到蛇怎么办?",DANGER),
       ("❓","迷路 — 我会怎么做?",CALM),
       ("🍱","食物 — 可以吃 vs 不可以吃",SAFE)]
for i,(em,cn,cl) in enumerate(pages):
    y=1.45+i*0.40
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(y),Inches(4.65),Inches(0.36))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(1.5)
    tb(s,5.05,y+0.04,0.5,0.30,em,sz=18,a=PP_ALIGN.CENTER)
    tb(s,5.55,y+0.04,4.0,0.30,cn,sz=13,b=True,c=cl)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.30))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=NAVY;sh3.line.width=Pt(2)
tb(s,0.5,3.95,9,0.35,"🗣️ 每页写 + 画  Each Page: Write + Draw",sz=14,b=True,c=NAVY)
tb(s,0.5,4.30,4.5,0.35,"·  画一个图 (Draw a picture)",sz=13,c=DARK)
tb(s,0.5,4.65,4.5,0.35,"·  写 1-2 句 (Write 1-2 sentences)",sz=13,c=DARK)
tb(s,5.2,4.30,4.5,0.35,"·  「我不 ___」 (Don't…)",sz=13,c=DARK)
tb(s,5.2,4.65,4.5,0.35,"·  「我会 ___」 (I will…)",sz=13,c=DARK)
n+=1;pn(s,n)
notes(s,"小书 / Poster (20-25 分钟):\n• 老师准备 4 页模板 (或用一张大纸折 4 折)\n• 每页一个主题\n• 学生自由设计: 画图 + 写 1-2 句\n• 完成后在教室展示 (墙上)\n• 句型固定: 「我不 ___」「我会 ___」")

# 34. ACTIVITY 3 — Emergency Whistle Workshop (bonus)
s=ns();bg(s,CREAM);hb(s,"🆘 ACTIVITY 3: 紧急哨子工坊  Whistle Workshop",SAFE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.40))
sh.fill.solid();sh.fill.fore_color.rgb=SAFE;sh.line.fill.background()
tb(s,0.4,0.99,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.0,"🎺 小哨子 (每人 1 个) — 老师准备",sz=13,c=DARK)
ap(tf,"🪢 绳子 / 项链 (挂哨子)",sz=13,c=DARK)
ap(tf,"🖍️ 彩笔 (装饰)",sz=13,c=DARK)
ap(tf,"📏 贴纸 (写名字)",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(0.40))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,4.95,0.99,4.65,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,4.95,1.45,4.7,1.6,"1️⃣ 装饰哨子 (彩笔 + 贴纸)",sz=13,c=DARK)
ap(tf2,"2️⃣ 写上名字 + 紧急联系",sz=13,c=DARK)
ap(tf2,"3️⃣ 系上绳子 → 挂脖子上",sz=13,c=DARK)
ap(tf2,"4️⃣ 练习: 紧急信号 = 3 长哨!",sz=13,b=True,c=DANGER)
ap(tf2,"5️⃣ 全班一起吹 3 长哨",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.55),Inches(9.4),Inches(1.55))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=DANGER;sh3.line.width=Pt(2.5)
tb(s,0.5,3.65,9.0,0.4,"⚠️ 国际求救信号  International SOS",sz=14,b=True,c=DANGER)
tb(s,0.5,4.05,9.0,0.5,"3 长哨 = 救我!  3 long whistles = HELP!",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,4.55,9.0,0.4,"等 30 秒 — 再 3 长哨 — 再等 — 一直重复",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"哨子工坊 (15-20 分钟):\n• 老师提前买 / 准备每人 1 个小哨子 ($1-2 一个)\n• 让学生装饰 (5-10 分钟)\n• 练习国际 SOS 信号: 3 长哨, 等 30 秒, 重复\n• 全班一起吹 (1 次就够了 — 太响)\n• 强调: 真的紧急才用!\n• 联系 Day 2 装包 (哨子是装包必备)")

# 35. MORE IDEAS — bonus
s=ns();bg(s,CREAM);hb(s,"💡 还可以做什么?  More Ideas",CAUTION)
tb(s,0.4,0.85,9.2,0.4,"如果时间充裕 — 试试这些加分活动!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
ideas=[
    ("🎯","危险 / 安全分类游戏","Sort: dangerous vs safe","老师拿出图片, 学生跑 「危险」/「安全」墙",CAUTION),
    ("🆘","紧急演练","Emergency drill","老师喊「迷路!」 学生快速走 6 步: 停-留-喊-哨-标-等",DANGER),
    ("📞","求救电话角色扮演","Pretend 911 call","2 人一组: 1 人求救, 1 人接听 — 练习报地址",CALM),
    ("🎨","「不能吃」海报","'Don't Eat' poster","小组做一张大海报 — 7 个不能碰物品",DANGER),
]
for i,(em,cn,en,desc,cl) in enumerate(ideas):
    y=1.40+i*0.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.78))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.55,y+0.20,0.7,0.4,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.05,4.0,0.35,cn,sz=15,b=True,c=cl)
    tb(s,1.30,y+0.40,4.0,0.30,en,sz=10,c=GRAY)
    tb(s,5.30,y+0.20,4.30,0.40,desc,sz=12,c=DARK)
n+=1;pn(s,n)

# 36. CLOSING — mission complete
s=ns();bg(s,DANGER)
tb(s,1,0.6,8,0.7,"🏆 任务完成!",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.35,8,0.5,"Mission Complete!",sz=20,c=WARM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.0),Inches(3.0),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=CAUTION;sh.line.color.rgb=WHITE;sh.line.width=Pt(4)
tb(s,3.5,2.4,3.0,0.6,"🆘",sz=80,a=PP_ALIGN.CENTER)
tb(s,3.5,3.6,3.0,0.5,"安全探险家",sz=22,b=True,c=DANGER,a=PP_ALIGN.CENTER)
tb(s,3.5,4.1,3.0,0.4,"Safety Explorer",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"恭喜! 8 天野外生存 — 你都学会了!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# === Save ===
out=os.path.join(BASE,"day4_emergency.pptx")
prs.save(out)
print(f"Saved {out}  ({n} slides)")
