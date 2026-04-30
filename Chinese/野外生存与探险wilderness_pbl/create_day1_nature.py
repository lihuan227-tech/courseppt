#!/usr/bin/env python3
"""
野外生存与探险 Wilderness Survival Unit — Day 1: 认识自然与安全规则
Structure modeled on 世界旅行 Day 1 Asia v2, adapted for 6 environments.
Palette: Adventure (pine green + sunset orange) — distinct from world-trip teal.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Adventure / Wilderness ---
PINE = RGBColor(0x1E,0x4D,0x2B)   # primary: deep pine green
SUN = RGBColor(0xE0,0x7A,0x2C)    # accent: sunset orange
CREAM = RGBColor(0xFD,0xF6,0xE3)  # background cream
BROWN = RGBColor(0x6B,0x44,0x23)  # soil brown
SKY = RGBColor(0x4A,0x90,0xD9)    # sky blue
SUNYEL = RGBColor(0xF5,0xC2,0x42) # sunshine yellow
ALERT = RGBColor(0xD0,0x4A,0x3C)  # warning red
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
WARM = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xE8,0xE8,0xE8)
GREEN_OK = RGBColor(0x38,0x8E,0x3C)

# Per-environment colors
FOREST = RGBColor(0x2D,0x5A,0x3D)
MOUNTAIN = RGBColor(0x54,0x6E,0x7A)
GRASS = RGBColor(0x7C,0xB3,0x42)
RIVER = RGBColor(0x19,0x76,0xD2)
DESERT = RGBColor(0xD4,0xA5,0x74)
SNOW = RGBColor(0x78,0xA7,0xB5)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=PINE,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color);tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER);tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER);return s
def vs(title,bgc):
    s=ns();bg(s,bgc);tb(s,1,0.8,8,0.8,"🎬 看视频  Watch Video",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,1.8,8,0.5,title,sz=22,c=WARM,a=PP_ALIGN.CENTER)
    ib(s,1.5,2.5,7,2.5,"📷 插入视频截图或粘贴视频链接");tb(s,1,5.1,8,0.3,"🔗 视频链接: ____________________",sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    return s

def aspect_slide(emoji,cn,en,env_color,aspect_label,aspect_color,questions,frame):
    """One slide for ONE aspect of ONE environment, inquiry-based (questions, no answers)."""
    s=ns();bg(s,CREAM)
    # Header bar (env color) — environment identity on left, aspect label on right
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.15),Inches(9.4),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=env_color;sh.line.fill.background()
    tb(s,0.5,0.22,1.0,0.55,emoji,sz=28,c=WHITE)
    tb(s,1.5,0.23,3.0,0.55,cn,sz=26,b=True,c=WHITE)
    tb(s,1.5,0.62,3.0,0.25,en,sz=11,c=WARM)
    # Aspect pill on the right of header
    pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.8),Inches(0.27),Inches(4.8),Inches(0.45))
    pill.fill.solid();pill.fill.fore_color.rgb=aspect_color;pill.line.color.rgb=WHITE;pill.line.width=Pt(1.5)
    tb(s,4.9,0.32,4.6,0.4,aspect_label,sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # Image placeholder (left)
    ib(s,0.3,1.05,4.3,3.3,f"📷 {cn} 图片 / 视频")
    # Inquiry panel (right) — questions, no answers
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(3.3))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE
    panel.line.color.rgb=aspect_color;panel.line.width=Pt(2.5)
    # Mini header inside the panel
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(0.5))
    head.fill.solid();head.fill.fore_color.rgb=aspect_color;head.line.fill.background()
    tb(s,5.0,1.13,4.6,0.4,"🤔 一起想一想  Let's Think Together",sz=14,b=True,c=WHITE)
    # Questions stacked with breathing room
    tf=tb(s,5.05,1.7,4.55,0.5,f"❓ {questions[0]}",sz=14,c=DARK)
    for q in questions[1:]:
        ap(tf,"",sz=8)
        ap(tf,f"❓ {q}",sz=14,c=DARK)
    # Sentence frame at the bottom (still inquiry — student fills the blanks)
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.5),Inches(9.4),Inches(0.65))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(2)
    tb(s,0.5,4.6,1.7,0.4,"💬 我来说",sz=14,b=True,c=SUN)
    tb(s,2.0,4.6,7.6,0.4,frame,sz=14,c=DARK)
    return s

def word_card_read(w,py,en,sent,img):
    """我会认 card — large character + pinyin + sentence."""
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",SUN)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=SUN,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=SUN;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=SUN)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,img):
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",PINE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=PINE)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=PINE)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    return s

# ============================================================
# UPGRADE HELPERS — Explorer Mission, video interaction, games, frames
# ============================================================
def notes(s,txt):
    s.notes_slide.notes_text_frame.text=txt

def mission_narrative_slide():
    s=ns();bg(s,CREAM);hb(s,"🧭 你是小探险家!  You're a Little Explorer!",PINE)
    hb_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(9.2),Inches(0.95))
    hb_box.fill.solid();hb_box.fill.fore_color.rgb=WARM;hb_box.line.color.rgb=SUN;hb_box.line.width=Pt(2.5)
    tb(s,0.6,1.10,8.8,0.45,"🌍 今天你要去 6 个地方探险!",sz=22,b=True,c=PINE,a=PP_ALIGN.CENTER)
    tb(s,0.6,1.55,8.8,0.35,"Today you visit 6 wild places — pack your imagination!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    envs=[("🌲","森林",FOREST),("🏔️","山地",MOUNTAIN),("🌾","草地",GRASS),
          ("🏞️","河边",RIVER),("🏜️","沙漠",DESERT),("❄️","雪地",SNOW)]
    for i,(em,cn,cl) in enumerate(envs):
        x=0.3+i*1.6;y=2.2
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(1.5),Inches(1.6))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
        tb(s,x+0.1,y+0.2,1.3,0.7,em,sz=34,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.95,1.3,0.4,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+1.3,1.3,0.25,f"#{i+1}",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.2),Inches(9.4),Inches(1.0))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(2)
    tb(s,0.5,4.30,1.7,0.4,"💬 我来说",sz=14,b=True,c=SUN)
    tb(s,2.0,4.27,7.6,0.3,"我是小探险家, 我要去 ____ 。",sz=15,b=True,c=DARK)
    tb(s,2.0,4.57,7.6,0.3,"I'm a little explorer. I'm going to ___.",sz=11,c=GRAY)
    tb(s,2.0,4.85,7.6,0.3,"举手做敬礼! Raise your hand and salute! 🫡",sz=12,b=True,c=PINE)
    return s

def mission_overview_slide():
    s=ns();bg(s,CREAM);hb(s,"🎒 6 个任务 · 6 颗 ⭐",SUN)
    tb(s,0.4,0.85,9.2,0.35,"每个地方一个任务 → 完成 1 个 = 拿 1 ⭐  ·  Earn 1 star per mission!",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    missions=[
        ("🌲","森林",FOREST,    "听 3 种声音",        "Listen for 3 sounds"),
        ("🏔️","山地",MOUNTAIN,  "决定: 上山 or 下山?","Decide: up or down?"),
        ("🌾","草地",GRASS,     "找草丛里的东西",     "Find what's hidden"),
        ("🏞️","河边",RIVER,     "安全 vs 危险?",      "Safe or dangerous?"),
        ("🏜️","沙漠",DESERT,    "选: 喝水 / 阴凉?",   "Choose: water or shade?"),
        ("❄️","雪地",SNOW,      "找最暖的地方",       "Find the warm spot"),
    ]
    for i,(em,cn,cl,m_cn,m_en) in enumerate(missions):
        col=i%3;row=i//3
        x=0.3+col*3.2;y=1.20+row*1.80
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.65))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
        tb(s,x+0.1,y+0.05,0.7,0.5,em,sz=26)
        tb(s,x+0.85,y+0.1,2.1,0.4,cn,sz=18,b=True,c=cl)
        tb(s,x+0.1,y+0.6,2.8,0.45,f"🎯 {m_cn}",sz=13,b=True,c=DARK)
        tb(s,x+0.1,y+1.05,2.8,0.35,m_en,sz=10,c=GRAY)
        tb(s,x+0.1,y+1.32,2.8,0.3,"→ 拿 1 ⭐",sz=12,b=True,c=SUN)
    # ⭐⭐⭐⭐⭐⭐ Star tracker bar
    tk=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.3),Inches(4.95),Inches(7.4),Inches(0.55))
    tk.fill.solid();tk.fill.fore_color.rgb=SUNYEL;tk.line.color.rgb=SUN;tk.line.width=Pt(2)
    tb(s,1.3,4.98,7.4,0.5,"☆  ☆  ☆  ☆  ☆  ☆   →   一共 6 颗 ⭐!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
    return s

def env_interaction_slide(em,cn,en_name,env_color,video_desc,before,during,after,
                          either_q,opt_a,opt_b,answer,reason,move_cn,move_en):
    s=ns();bg(s,CREAM);hb(s,f"{em} {cn} · 🎬 探险任务  Mission  ·  完成 = ⭐",env_color)
    # Thin video strip (top)
    v=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.85),Inches(9.4),Inches(0.45))
    v.fill.solid();v.fill.fore_color.rgb=env_color;v.line.fill.background()
    tb(s,0.45,0.91,1.6,0.35,"🎬 视频",sz=12,b=True,c=WHITE)
    tb(s,2.0,0.91,7.6,0.35,video_desc,sz=12,b=True,c=WHITE)
    # B / D / A — three cards across
    bda=[("👀 看之前 Before",before),
         ("🔎 看时 During",during),
         ("💬 看之后 After",after)]
    for i,(lbl,content) in enumerate(bda):
        x=0.3+i*3.15
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(3.05),Inches(1.05))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=env_color;sh.line.width=Pt(2)
        head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(3.05),Inches(0.35))
        head.fill.solid();head.fill.fore_color.rgb=env_color;head.line.fill.background()
        tb(s,x+0.1,1.48,2.9,0.3,lbl,sz=11,b=True,c=WHITE)
        tb(s,x+0.1,1.85,2.9,0.6,content,sz=11,c=DARK)
    # Either/or (left, smaller)
    eo=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(2.65),Inches(6.4),Inches(2.10))
    eo.fill.solid();eo.fill.fore_color.rgb=WARM;eo.line.color.rgb=ALERT;eo.line.width=Pt(2)
    pb2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.45),Inches(2.74),Inches(2.4),Inches(0.4))
    pb2.fill.solid();pb2.fill.fore_color.rgb=ALERT;pb2.line.fill.background()
    tb(s,0.55,2.79,2.3,0.35,"🤔 你来选 Choose",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,2.95,2.79,3.6,0.35,either_q,sz=11,b=True,c=DARK)
    a_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.45),Inches(3.22),Inches(3.0),Inches(1.45))
    a_box.fill.solid();a_box.fill.fore_color.rgb=WHITE;a_box.line.color.rgb=PINE;a_box.line.width=Pt(2)
    tb(s,0.55,3.27,2.8,0.4,"A",sz=14,b=True,c=PINE)
    tb(s,0.55,3.62,2.8,1.0,opt_a,sz=11,c=DARK)
    b_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.55),Inches(3.22),Inches(3.05),Inches(1.45))
    b_box.fill.solid();b_box.fill.fore_color.rgb=WHITE;b_box.line.color.rgb=ALERT;b_box.line.width=Pt(2)
    tb(s,3.65,3.27,2.85,0.4,"B",sz=14,b=True,c=ALERT)
    tb(s,3.65,3.62,2.85,1.0,opt_b,sz=11,c=DARK)
    # Movement (right)
    mv=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.95),Inches(2.65),Inches(2.75),Inches(2.10))
    mv.fill.solid();mv.fill.fore_color.rgb=GREEN_OK;mv.line.fill.background()
    tb(s,7.05,2.74,2.55,0.35,"🚶 演一演 Act!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,7.05,3.18,2.55,0.5,move_cn,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,7.05,3.72,2.55,0.35,move_en,sz=10,c=WARM,a=PP_ALIGN.CENTER)
    tb(s,7.05,4.18,2.55,0.3,"全班 30 秒",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,7.05,4.42,2.55,0.25,"Whole class · 30s",sz=9,c=WARM,a=PP_ALIGN.CENTER)
    # ⭐ + sentence frame
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.85),Inches(9.4),Inches(0.4))
    sf.fill.solid();sf.fill.fore_color.rgb=SUNYEL;sf.line.color.rgb=SUN;sf.line.width=Pt(1.5)
    tb(s,0.45,4.91,9.0,0.3,"💬 我选 ___, 因为 ___ 。   →   完成挑战 = 拿 1 ⭐",sz=12,b=True,c=DARK)
    notes(s,f"老师备课:\n• 答案 Answer: {answer}\n• 原因 Reason: {reason}\n• 视频建议: {video_desc}\n• 全班动作 30 秒 — 站起来一起做。\n• A/B: 把教室分两边, 学生站到选择的一边, 说出原因。\n• 完成 = 全班拿 1 ⭐ → 涂在白板上的星图。")
    return s

def gesture_game_slide():
    s=ns();bg(s,CREAM);hb(s,"🛡️ 安全 vs 危险!  Safe vs Dangerous!",ALERT)
    # Rules box
    rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.0),Inches(9.4),Inches(1.1))
    rb.fill.solid();rb.fill.fore_color.rgb=WARM;rb.line.color.rgb=ALERT;rb.line.width=Pt(2)
    tb(s,0.5,1.10,9.0,0.4,"老师说一件事 — 学生用动作回答!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.5,1.50,4.4,0.5,"✅ 安全 = 双手举高 (V)",sz=15,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
    tb(s,5.1,1.50,4.4,0.5,"❌ 危险 = 双手交叉 (X)",sz=15,b=True,c=ALERT,a=PP_ALIGN.CENTER)
    # 6 examples
    examples=[("🌲","摸不认识的蘑菇",       "❌"),
              ("🏔️","在山顶看天有黑云就下山","✅"),
              ("🌾","看到蜜蜂就站着不动",   "✅"),
              ("🏞️","河边的水可以直接喝",   "❌"),
              ("🏜️","太阳下走 1 小时不喝水","❌"),
              ("❄️","看到河面是冰就走过去", "❌")]
    for i,(em,desc,ans) in enumerate(examples):
        col=i%2;row=i//2
        x=0.3+col*4.7;y=2.3+row*0.85
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(0.75))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ALERT;sh.line.width=Pt(1.5)
        tb(s,x+0.1,y+0.15,0.6,0.5,em,sz=22)
        tb(s,x+0.75,y+0.18,3.2,0.4,desc,sz=12,b=True,c=DARK)
        tb(s,x+3.95,y+0.15,0.5,0.5,ans,sz=22,a=PP_ALIGN.CENTER)
    notes(s,"老师玩法 (3-5 分钟):\n• 念上面 6 个情境, 学生举手或交叉手回答。\n• 答错的不出局, 让他们再听一次。\n• G1-3 学生加一句: 「___ 不安全, 因为 ___」")
    return s

def where_am_i_slide():
    s=ns();bg(s,CREAM);hb(s,"🔍 我在哪里？  Where Am I?",SUN)
    tb(s,0.4,0.95,9.2,0.4,"老师演动作或说一个声音 — 学生猜环境!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.35,9.2,0.3,"Teacher acts/sounds — students guess the place!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    clues=[("🌲","老师弯腰拨树枝、做听虫子的动作","Bending, pushing branches",FOREST),
           ("🏔️","用手扶岩壁、深呼吸",         "Climbing, deep breath",MOUNTAIN),
           ("🌾","眯眼看远方、伸手摸花",       "Squinting, touching flowers",GRASS),
           ("🏞️","「哗哗」水声、小心踩石头",  "Water sounds, stepping",RIVER),
           ("🏜️","擦汗、喝水、慢慢走",       "Wiping sweat, drinking",DESERT),
           ("❄️","发抖、跺脚、抱身体",       "Shivering, stomping",SNOW)]
    for i,(em,cn,en,cl) in enumerate(clues):
        col=i%3;row=i//3
        x=0.3+col*3.2;y=1.85+row*1.55
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.4))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
        tb(s,x+0.1,y+0.1,2.8,0.5,em,sz=24,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.6,2.8,0.4,cn,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.98,2.8,0.35,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.95),Inches(9.4),Inches(0.4))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(1.5)
    tb(s,0.5,5.0,9.0,0.3,"💬 你在 ____ 。/ I'm in ___.    G1-3: 我猜你在 ___, 因为 ___ 。",sz=12,b=True,c=DARK)
    notes(s,"低 prep — 老师当场演 / 发声音, 学生举手猜。每次 3-5 个情境。")
    return s

def sentence_frames_slide():
    s=ns();bg(s,CREAM);hb(s,"💬 句型卡  Sentence Frames (K · G1–3)",SUN)
    # K column — 5 frames
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.55),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
    pb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.45),Inches(1.05),Inches(1.6),Inches(0.4))
    pb.fill.solid();pb.fill.fore_color.rgb=SUN;pb.line.fill.background()
    tb(s,0.55,1.10,1.5,0.35,"K (TK-K)",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    k_frames=[("这是 ____ 。",      "This is ___."),
              ("很 ____ 。",        "Very ___."),
              ("我看到 ____ 。",    "I see ___."),
              ("安全 / 不安全",     "Safe / not safe"),
              ("可以 / 不可以",     "Can / cannot")]
    for i,(cn,en) in enumerate(k_frames):
        y=1.55+i*0.62
        tb(s,0.5,y,4.3,0.4,f"·  {cn}",sz=18,b=True,c=PINE)
        tb(s,0.5,y+0.32,4.3,0.25,en,sz=9,c=GRAY)
    # G1-3 column — 4 frames
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.95),Inches(4.65),Inches(4.0))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=PINE;sh2.line.width=Pt(2.5)
    pb2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.20),Inches(1.05),Inches(1.7),Inches(0.4))
    pb2.fill.solid();pb2.fill.fore_color.rgb=PINE;pb2.line.fill.background()
    tb(s,5.30,1.10,1.6,0.35,"G1 - G3",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    g_frames=[("这里很 ___, 因为 ___ 。",     "It's very ___ here, because ___."),
              ("我应该 ___ 。",                "I should ___."),
              ("___ 不安全, 因为 ___ 。",      "___ is not safe, because ___."),
              ("我选 ___, 因为 ___ 。",        "I choose ___, because ___.")]
    for i,(cn,en) in enumerate(g_frames):
        y=1.65+i*0.78
        tb(s,5.25,y,4.4,0.45,f"·  {cn}",sz=16,b=True,c=PINE)
        tb(s,5.25,y+0.4,4.4,0.25,en,sz=9,c=GRAY)
    # Bottom hint
    tb(s,0.4,5.1,9.2,0.3,"💡 把这张 PPT 截屏打印, 贴在每张桌子上 — 学生整堂课参考。",sz=11,b=True,c=BROWN,a=PP_ALIGN.CENTER)
    notes(s,"零 prep: 打印这张当桌签。\nK 重点: 2-字短语 + 短句 (这是, 很, 我看到, 安全, 可以)。\nG1-3: 用 因为 + 应该 + 选 串成完整句。")
    return s

# --- New 2-slide-per-env helpers — 5-stage flow (看·想·说·判·做) ---

def experience_slide(em,cn,en_name,env_color,see_q,hear_q,feel_q,think_chips):
    """SLIDE 1 of 2 — Stage 1️⃣ 看 (Observe) + Stage 2️⃣ 想 (Think).
    think_chips: list of 4 (emoji, short_q) tuples for brainstorming hints."""
    s=ns();bg(s,CREAM);hb(s,f"{em} {cn}  {en_name} · 1️⃣ 看 + 2️⃣ 想  Observe + Think",env_color)
    # === 1️⃣ 看 (Observe) — top half ===
    # Image (left)
    ib(s,0.3,0.95,4.4,2.20,f"📷 {cn} 图片 / 视频")
    # Sensory panel (right)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(2.20))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE
    panel.line.color.rgb=env_color;panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(0.40))
    head.fill.solid();head.fill.fore_color.rgb=env_color;head.line.fill.background()
    tb(s,5.0,0.99,4.6,0.35,"1️⃣ 👀 看 — 你看到 / 听到 / 感觉到什么?",sz=12,b=True,c=WHITE)
    # 3 sensory cards (compact)
    rows=[("👀",see_q),("👂",hear_q),("✋",feel_q)]
    for i,(icon,q) in enumerate(rows):
        y=1.45+i*0.55
        tb(s,5.05,y,0.5,0.45,icon,sz=22)
        tb(s,5.65,y+0.05,4.0,0.45,q,sz=13,b=True,c=DARK)
    # === 2️⃣ 想 (Think) — bottom-half panel with hint chips ===
    tp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.30),Inches(9.4),Inches(1.30))
    tp.fill.solid();tp.fill.fore_color.rgb=WARM
    tp.line.color.rgb=SUN;tp.line.width=Pt(2)
    th=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.30),Inches(9.4),Inches(0.42))
    th.fill.solid();th.fill.fore_color.rgb=SUN;th.line.fill.background()
    tb(s,0.45,3.34,9.2,0.35,"2️⃣ 🤔 想 — 如果你去那里, 会发生什么?  What might happen?",sz=12,b=True,c=WHITE)
    # 4 hint chips
    for i,(chip_em,chip_t) in enumerate(think_chips):
        x=0.45+i*2.30
        ch=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.85),Inches(2.20),Inches(0.65))
        ch.fill.solid();ch.fill.fore_color.rgb=WHITE
        ch.line.color.rgb=SUN;ch.line.width=Pt(1.5)
        tb(s,x+0.1,3.92,0.6,0.5,chip_em,sz=22)
        tb(s,x+0.75,3.95,1.40,0.55,chip_t,sz=12,b=True,c=DARK)
    # Frames bar at bottom (K + G2-G3)
    fy=4.70
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(fy),Inches(9.4),Inches(0.55))
    sf.fill.solid();sf.fill.fore_color.rgb=SUNYEL
    sf.line.color.rgb=SUN;sf.line.width=Pt(1.5)
    tb(s,0.5,fy+0.03,1.3,0.25,"💬 K-G1:",sz=10,b=True,c=ALERT)
    tb(s,1.7,fy+0.03,7.9,0.25,"我看到 ___ 。",sz=11,b=True,c=DARK)
    tb(s,0.5,fy+0.30,1.3,0.25,"💬 G2-G3:",sz=10,b=True,c=PINE)
    tb(s,1.7,fy+0.30,7.9,0.25,"我觉得会 ___ 。",sz=11,b=True,c=DARK)
    return s

def safety_slide(em,cn,en_name,env_color,say_summary,situations,move_cn):
    """SLIDE 2 of 2 — Stage 3️⃣ 说 + 4️⃣ 判 + 5️⃣ 做 (Say + Decide + Act).
    say_summary: short teacher-summary line on key situations
    situations: list of 3 (question, opt_a, opt_b, correct_idx)
    move_cn: physical action prompt"""
    s=ns();bg(s,CREAM);hb(s,f"{em} {cn} · 3️⃣ 说 + 4️⃣ 判 + 5️⃣ 做  Say · Decide · Act  ⭐",ALERT)
    # === 3️⃣ 说 (Say) — teacher summary band ===
    sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.92),Inches(9.4),Inches(0.45))
    sb.fill.solid();sb.fill.fore_color.rgb=BROWN;sb.line.fill.background()
    tb(s,0.45,0.96,1.7,0.35,"3️⃣ 📢 说",sz=12,b=True,c=SUNYEL)
    tb(s,2.0,0.96,7.6,0.35,say_summary,sz=12,b=True,c=WHITE)
    # === 4️⃣ 判 (Decide) — 3 compact A/B rows ===
    sit_h=0.74
    base_y=1.75
    # Section header line (between 说 band and rows)
    tb(s,0.3,1.42,9.4,0.25,"4️⃣ ⚖️ 判 — 选 A 还是 B?  Choose A or B?",sz=11,b=True,c=ALERT)
    for i,(q,opt_a,opt_b,correct) in enumerate(situations):
        y=base_y+i*(sit_h+0.05)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(y),Inches(9.4),Inches(sit_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE
        sh.line.color.rgb=ALERT;sh.line.width=Pt(1.5)
        nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.40),Inches(y+0.06),Inches(0.32),Inches(0.32))
        nb.fill.solid();nb.fill.fore_color.rgb=ALERT;nb.line.fill.background()
        tb(s,0.40,y+0.08,0.32,0.28,str(i+1),sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,0.85,y+0.05,8.5,0.30,q,sz=12,b=True,c=DARK)
        # A and B side by side
        a_correct=(correct==0);b_correct=(correct==1)
        a_color=GREEN_OK if a_correct else ALERT
        b_color=GREEN_OK if b_correct else ALERT
        a_fill=WARM if a_correct else WHITE
        b_fill=WARM if b_correct else WHITE
        a_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y+0.40),Inches(4.3),Inches(0.32))
        a_box.fill.solid();a_box.fill.fore_color.rgb=a_fill
        a_box.line.color.rgb=a_color;a_box.line.width=Pt(1.2)
        tb(s,0.6,y+0.42,0.4,0.3,"A.",sz=11,b=True,c=a_color)
        tb(s,1.0,y+0.42,3.7,0.3,opt_a,sz=11,c=DARK)
        b_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(y+0.40),Inches(4.3),Inches(0.32))
        b_box.fill.solid();b_box.fill.fore_color.rgb=b_fill
        b_box.line.color.rgb=b_color;b_box.line.width=Pt(1.2)
        tb(s,5.3,y+0.42,0.4,0.3,"B.",sz=11,b=True,c=b_color)
        tb(s,5.7,y+0.42,3.7,0.3,opt_b,sz=11,c=DARK)
    # === 5️⃣ 做 (Act) — movement strip ===
    act_y=base_y+3*(sit_h+0.05)+0.05
    mv=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(act_y),Inches(9.4),Inches(0.40))
    mv.fill.solid();mv.fill.fore_color.rgb=GREEN_OK;mv.line.fill.background()
    tb(s,0.45,act_y+0.04,1.7,0.32,"5️⃣ 🚶 做 30s:",sz=12,b=True,c=SUNYEL)
    tb(s,2.0,act_y+0.04,7.6,0.32,move_cn,sz=12,b=True,c=WHITE)
    # === Frames bar at bottom ===
    fy=act_y+0.50
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(fy),Inches(9.4),Inches(0.50))
    sf.fill.solid();sf.fill.fore_color.rgb=SUNYEL
    sf.line.color.rgb=SUN;sf.line.width=Pt(1.5)
    tb(s,0.5,fy+0.03,1.3,0.25,"💬 K-G1:",sz=10,b=True,c=ALERT)
    tb(s,1.7,fy+0.03,7.9,0.25,"不可以 ___ 。 安全 / 不安全。",sz=11,b=True,c=DARK)
    tb(s,0.5,fy+0.27,1.3,0.25,"💬 G2-G3:",sz=10,b=True,c=PINE)
    tb(s,1.7,fy+0.27,7.9,0.25,"___ 不安全, 因为 ___ 。 我应该 ___ 。",sz=11,b=True,c=DARK)
    return s

n=0

# ============================================================
# 1 COVER — Explorer Badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.25,8,0.7,"Wilderness Adventure Camp",sz=32,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,1,0.85,8,0.45,"野外生存与探险夏令营",sz=20,c=PINE,a=PP_ALIGN.CENTER)
# Big round badge
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.5),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=PINE;sh.line.color.rgb=SUN;sh.line.width=Pt(6)
# Inner badge ring
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.8),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=SUN;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.0,2.8,0.4,"DAY 1",sz=16,b=True,c=SUN,a=PP_ALIGN.CENTER)
ap(tf,"🏕️",sz=60,a=PP_ALIGN.CENTER)
ap(tf,"认识自然",sz=20,b=True,c=PINE,a=PP_ALIGN.CENTER)
ap(tf,"NATURE & SAFETY",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,5.05,8,0.4,"🎒 准备好背包，我们出发！Let's go, explorers!",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","认识6种自然环境 + 户外安全规则",PINE),
    ("Session 2  下午","2:00-2:45","复习总结 + 语言目标 (认字写字)",SUN),
    ("Session 3  下午","3:00-4:30","写Booklet + Project 项目活动",BROWN)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=WARM)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"🌲 内容目标  Content:",sz=20,b=True,c=PINE)
tf=tb(s,0.7,1.4,9,1.3,"1. 了解 6 种常见自然环境：森林、山地、草地、河边、沙漠、雪地",sz=15,c=DARK)
ap(tf,"2. 了解不同自然环境的基本特点",sz=15,c=DARK)
ap(tf,"3. 了解户外可能遇到的危险与基本安全规则",sz=15,c=DARK)
ap(tf,"4. 建立基本的户外安全意识",sz=15,c=DARK)
tb(s,0.5,3.0,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=SUN)
tb(s,0.7,3.5,4.4,0.9,"👀 我会认：森林 山地 草地\n　　　　　河边 沙漠 雪地",sz=15,b=True,c=DARK)
tb(s,5.3,3.5,4.3,0.9,"✍️ 我会写：森林 山地 河边",sz=15,b=True,c=DARK)
tb(s,0.5,4.55,9,0.5,"🎨 实践目标：完成 Booklet + 2 个项目 + 1 个游戏",sz=15,c=BROWN)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","认识 6 种自然环境 + 安全规则\n🌲 森林  🏔️ 山地  🌾 草地  🏞️ 河边  🏜️ 沙漠  ❄️ 雪地",PINE,"🧭")
n+=1

# ============================================================
# 4-NEW  EXPLORER MISSION (narrative + 6 missions)
# ============================================================
s=mission_narrative_slide();n+=1;pn(s,n)
# s=mission_overview_slide();n+=1;pn(s,n)  # removed per request — no more star tracker overview

# ============================================================
# 4a PICTURE BOOK INTRO — Story Time hook
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📖 先来听个故事  Story Time",SUN)
# Left: book illustration placeholder
ib(s,0.3,1.0,4.4,3.6,"📷 绘本封面 / Book Cover")
# Right: how to watch + link
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.0),Inches(4.8),Inches(3.6))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.0),Inches(4.8),Inches(0.5))
head.fill.solid();head.fill.fore_color.rgb=PINE;head.line.fill.background()
tb(s,5.05,1.08,4.6,0.4,"🎬 一起看绘本  Watch Together",sz=15,b=True,c=WHITE)
tf=tb(s,5.1,1.7,4.55,0.4,"👂 听一听：故事里的人去了哪里？",sz=14,c=DARK)
ap(tf,"",sz=6)
ap(tf,"👀 看一看：他们看到了什么？",sz=14,c=DARK)
ap(tf,"",sz=6)
ap(tf,"🤔 想一想：哪里最危险？",sz=14,c=DARK)
ap(tf,"",sz=10)
ap(tf,"🔗 绘本链接 Link:",sz=12,b=True,c=SUN)
ap(tf,"limaogushi.com/play?id=783",sz=12,c=SKY)
# Bottom: bridge to discussion
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.75),Inches(9.4),Inches(0.5))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(2)
tb(s,0.5,4.83,9.0,0.4,"📌 看完后，我们一起聊一聊故事里的地方！",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 4b DISCUSSION — Questions after the book + sentence frames
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"❓ 故事过后  Let's Discuss",SUN)
# Two main question cards
discussion_qs=[
    ("1️⃣","他们去了哪些地方？","Where did they go?",PINE),
    ("2️⃣","这些地方有什么特点？可能有什么危险？","Features and dangers of these places?",ALERT),
]
for i,(num,q_cn,q_en,cl) in enumerate(discussion_qs):
    y=1.0+i*1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.6),Inches(y+0.18),Inches(0.65),Inches(0.65))
    nb.fill.solid();nb.fill.fore_color.rgb=cl;nb.line.fill.background()
    tb(s,0.6,y+0.22,0.65,0.55,num,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.5,y+0.13,8.0,0.5,q_cn,sz=19,b=True,c=DARK)
    tb(s,1.5,y+0.6,8.0,0.35,q_en,sz=12,c=GRAY)
# Sentence frame panel — students answer using these patterns, can extend beyond the book
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.4),Inches(9.2),Inches(1.85))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.4),Inches(9.2),Inches(0.45))
head.fill.solid();head.fill.fore_color.rgb=SUN;head.line.fill.background()
tb(s,0.6,3.45,9.0,0.4,"💬 我会说  Sentence Frames (绘本以外的地方也欢迎说！)",sz=14,b=True,c=WHITE)
# 4 frames in 2x2 grid
frames=[
    "他们去了 __________。",
    "这里很 __________。",
    "这里可能有 __________ 危险。",
    "我们要 __________。",
]
for i,fr in enumerate(frames):
    col=i%2;row=i//2
    x=0.6+col*4.55;y=3.95+row*0.55
    tb(s,x,y,4.4,0.4,f"·  {fr}",sz=14,c=DARK)
pn(s,n)

# ============================================================
# 4c TEACHER NOTE — Discussion guide table
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📋 Teacher Note  课堂提问参考",BROWN)
tb(s,0.4,0.9,9.2,0.35,"引导孩子说出每个地方的「特点」和「可能的危险」。鼓励孩子拓展绘本以外自己知道的内容。",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Table: 7 rows (header + 6), 3 columns
ts=s.shapes.add_table(7,3,Inches(0.3),Inches(1.3),Inches(9.4),Inches(3.5));t=ts.table
t.columns[0].width=Inches(1.5)
t.columns[1].width=Inches(3.4)
t.columns[2].width=Inches(4.5)
table_rows=[
    ["地方  Place","特点  Features","可能的危险  Possible Dangers"],
    ["🌾 草地","草很高，看不清脚下","可能摔倒、被虫咬、踩到石头或坑"],
    ["🏞️ 河边 / 河流","有水，地面湿滑","可能滑倒、掉进水里、水流太急"],
    ["🟫 泥地","又湿又软，走路困难","可能滑倒、鞋子陷进去"],
    ["🌲 森林","树很多，光线暗","容易迷路，可能有不认识的植物或动物"],
    ["❄️ 雪地","很冷，地面滑","可能滑倒、着凉、看不清路"],
    ["🕳️ 山洞","黑黑的，看不清里面","可能有野生动物，不安全，不能随便进去"],
]
for r,rd in enumerate(table_rows):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True
        p=tf.paragraphs[0];p.alignment=PP_ALIGN.LEFT if r>0 else PP_ALIGN.CENTER
        rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        if r==0:
            rn.font.size=Pt(13);rn.font.bold=True;rn.font.color.rgb=WHITE
            cl.fill.solid();cl.fill.fore_color.rgb=BROWN
        else:
            rn.font.size=Pt(12);rn.font.color.rgb=DARK
            if c==0:rn.font.bold=True
            if r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF0,0xE8)
            else:cl.fill.solid();cl.fill.fore_color.rgb=WHITE
# Tip at bottom
tip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.9),Inches(9.4),Inches(0.5))
tip.fill.solid();tip.fill.fore_color.rgb=WARM;tip.line.color.rgb=SUN;tip.line.width=Pt(2)
tb(s,0.5,4.97,9.0,0.4,"💡 让孩子拓展：除了绘本里的地方，你还知道哪里有危险？(操场？停车场？)",sz=12,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 5 OVERVIEW — 6 environments preview
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 大自然里有什么？  What's in Nature?")
tb(s,0.4,0.9,9,0.4,"今天我们一起认识 6 种常见自然环境！",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
envs=[
    ("🌲","森林","Forest",FOREST),
    ("🏔️","山地","Mountain",MOUNTAIN),
    ("🌾","草地","Grassland",GRASS),
    ("🏞️","河边","Riverside",RIVER),
    ("🏜️","沙漠","Desert",DESERT),
    ("❄️","雪地","Snow",SNOW),
]
for i,(em,cn,en,cl) in enumerate(envs):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.45+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.1,2.8,0.7,em,sz=40,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.95,2.8,0.45,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.4,2.8,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# Sentence frames moved INTO each env's Experience + Safety slide (per-env, leveled)

# ============================================================
# 6-29  Inquiry slides — 6 environments × 4 aspects = 24 slides
#       Each aspect gets its own slide. Students discover answers
#       through questions (inquiry-based), not through statements.
# ============================================================
ASPECTS=[
    ("👀 长什么样  Looks Like", FOREST,   "我看到 ____ 和 ____。这里 ____。"),
    ("✨ 特点  Features",       SUN,      "____ 里有 ____。这里的 ____ 让我 ____。"),
    ("⚠️ 小心危险  Dangers",    ALERT,    "在这里要小心 ____ 和 ____。"),
    ("✅ 安全规则  Safety Rules",GREEN_OK, "我应该 ____。我不能 ____。"),
]

inquiry_data=[
    ("🌲","森林","Forest",FOREST,[
        # Looks Like (observation — kept concrete)
        ["森林的光线 — 亮 还是 暗?",
         "树是 多 还是 少? 高 还是 矮?",
         "你最先看到什么?"],
        # Features
        ["森林里 — 冷 还是 暖?",
         "空气闻起来 — 香 还是 没味道?",
         "森林里有什么动物?"],
        # Dangers (K-5 either/or)
        ["看到不认识的蘑菇 — 能吃 还是 不能吃?",
         "听到「沙沙」声 — 大喊 还是 安静?",
         "看到野生动物 — 走近 还是 站住?"],
        # Safety Rules (K-5 action-based)
        ["迷路了 — 跑 还是 留在原地?",
         "不认识的果子 — 尝一口 还是 不要碰?",
         "走森林 — 一个人 还是 跟着大人?"],
    ]),
    ("🏔️","山地","Mountain",MOUNTAIN,[
        ["山 — 高 还是 矮?",
         "山顶 — 能看远 还是 看不远?",
         "山上 — 树多 还是 石头多?"],
        ["山上和山下 — 哪边冷?",
         "山上的风 — 大 还是 小?",
         "山上很安静 — 你能听到什么?"],
        ["山顶很滑 — 走快 还是 走慢?",
         "看到大石头掉下来 — 摸 还是 不摸?",
         "天上有黑云 — 上山 还是 下山?"],
        ["爬山穿 — 拖鞋 还是 运动鞋?",
         "山上 — 冷 还是 热?",
         "下山时 — 手放口袋 还是 扶东西?"],
    ]),
    ("🌾","草地","Grassland",GRASS,[
        ["草地 — 平 还是 不平?",
         "草地是什么颜色?",
         "除了草, 你看到什么?"],
        ["太阳 — 晒 还是 不晒?",
         "草地里的虫 — 多 还是 少?",
         "草地 — 软 还是 硬?"],
        ["蜜蜂飞过来 — 挥手 还是 站着不动?",
         "草地走路 — 看脚下 还是 看天?",
         "太阳很大 — 戴帽子 还是 不戴?"],
        ["草地穿 — 短裤 还是 长裤?",
         "皮肤要 — 涂防晒 还是 不涂?",
         "看到小虫 — 拍打 还是 不动?"],
    ]),
    ("🏞️","河边","Riverside",RIVER,[
        ["你看到的水 — 在动 还是 不动?",
         "岸边 — 泥 还是 石头?",
         "水 — 干净 还是 浑浊?"],
        ["河水 — 冷 还是 热?",
         "水里 — 有鱼 还是 没鱼?",
         "河边树 — 多 还是 少?"],
        ["河水看起来浅 — 安全 还是 不安全?",
         "湿石头 — 跳过去 还是 慢慢走?",
         "想喝水 — 河水 还是 水壶水?"],
        ["下水玩 — 可以 还是 不可以?",
         "穿什么? — 拖鞋 还是 防滑鞋?",
         "身边要有 — 大人 还是 自己一个人?"],
    ]),
    ("🏜️","沙漠","Desert",DESERT,[
        ["沙漠 — 树多 还是 树少?",
         "沙漠 — 很热 还是 很冷?",
         "你看到什么? 沙、还是石头?"],
        ["沙漠里植物 — 高 还是 矮?",
         "沙漠的动物 — 大 还是 小?",
         "沙丘 — 圆 还是 平?"],
        ["中午太阳大 — 走 还是 找阴凉?",
         "一个人离队 — 可以 还是 不可以?",
         "渴了 — 喝水 还是 等等?"],
        ["穿 — 长袖 还是 短袖?",
         "戴 — 帽子 还是 不戴?",
         "水带 — 多 还是 少?"],
    ]),
    ("❄️","雪地","Snow",SNOW,[
        ["雪 — 软 还是 硬?",
         "雪是什么颜色? 冷不冷?",
         "踩雪的声音 — 你听过吗?"],
        ["空气 — 暖和 还是 寒冷?",
         "雪 — 软 还是 硬?",
         "雪地里能玩什么?"],
        ["雪地很冷 — 手套要 还是 不要?",
         "河面是冰 — 走过去 还是 绕道?",
         "太阳 + 雪 — 戴墨镜 还是 不戴?"],
        ["穿 — 厚衣 还是 短袖?",
         "走湖面 — 可以 还是 不可以?",
         "脚 — 穿雪靴 还是 拖鞋?"],
    ]),
]

# Per-env interaction packs (B / D / A + either-or + movement)
ENV_INTERACTIONS={
    "森林":dict(video_desc="森林里的声音 / Forest soundscape (1 min)",
        before="闭上眼睛\nClose your eyes",
        during="听 3 个声音\nListen for 3 sounds",
        after="模仿 1 个声音\nMimic 1 sound!",
        either_q="听到树丛「沙沙」响 — 怎么办?",
        opt_a="A. 跑开\nRun away",
        opt_b="B. 安静站着\nStand still & quiet",
        answer="B 安静站着",
        reason="跑会惊到野生动物, 它可能追过来。安静观察。",
        move_cn="弯腰拨树枝走路",move_en="Bend, push branches"),
    "山地":dict(video_desc="登山者爬山顶 / Climbers reaching summit (30s)",
        before="看登山者的脚\nWatch their feet",
        during="数 5 个小心的脚步\nCount 5 careful steps",
        after="学他爬山 (慢慢)\nAct slow climbing",
        either_q="山顶上有黑云 — 上还是下?",
        opt_a="A. 继续上\nKeep going up",
        opt_b="B. 下山\nGo back down",
        answer="B 下山",
        reason="山上的雷雨很危险, 闪电先打高的地方。",
        move_cn="扶着假岩壁, 一步一步爬",move_en="Climb, hold the wall"),
    "草地":dict(video_desc="草地虫子和花 / Insects and flowers (30s)",
        before="看草地什么颜色\nWhat color is it?",
        during="找 2 个动物\nFind 2 animals",
        after="模仿 1 个动物 1 秒\nMimic 1 animal · 1s",
        either_q="蜜蜂飞过来 — 怎么办?",
        opt_a="A. 挥手赶走\nWave hands",
        opt_b="B. 站着不动\nStand still",
        answer="B 站着不动",
        reason="挥手会让蜜蜂以为你要打它, 就会蛰你。慢慢走开。",
        move_cn="慢慢走 + 看脚下 + 戴帽子",move_en="Walk slow, look down"),
    "河边":dict(video_desc="孩子在河边玩 / Kids by river (30s)",
        before="看孩子离水多近\nHow close to water?",
        during="找 3 块湿石头\nSpot 3 wet rocks",
        after="假装小心走\nStep carefully",
        either_q="想喝水 — 哪个水?",
        opt_a="A. 河里的水\nRiver water",
        opt_b="B. 水壶里的水\nFrom the bottle",
        answer="B 水壶水",
        reason="河水可能有细菌、虫子 — 喝了会拉肚子。煮过的水才能喝。",
        move_cn="小心踩湿石头, 双手张开",move_en="Step on wet rocks carefully"),
    "沙漠":dict(video_desc="沙漠里的骆驼和探险家 / Camel + explorer (30s)",
        before="看探险家穿什么\nWhat is he wearing?",
        during="数他喝几口水\nCount his sips",
        after="擦汗 + 喝水\nWipe + drink!",
        either_q="正午太阳很大 — 怎么办?",
        opt_a="A. 多走快点\nWalk faster",
        opt_b="B. 找阴凉休息\nRest in the shade",
        answer="B 阴凉休息",
        reason="正午温度可达 50°C — 走快会中暑。等下午凉了再走。",
        move_cn="擦汗, 慢慢走, 喝水",move_en="Wipe, walk slow, drink"),
    "雪地":dict(video_desc="雪山上的探险家 / Snow mountain explorer (30s)",
        before="数他穿几件衣服\nCount layers",
        during="找他的手套\nFind his gloves",
        after="发抖 + 抱紧\nShiver + hug!",
        either_q="看到河面是冰 — 走过去?",
        opt_a="A. 走过去\nWalk across",
        opt_b="B. 绕道走\nGo around",
        answer="B 绕道走",
        reason="冰可能是薄的, 看不出来。掉进冰水里非常危险。",
        move_cn="发抖, 跺脚, 抱紧自己",move_en="Shiver, stomp, hug self"),
}

# NEW STRUCTURE: 2 slides per environment
# Slide 1: EXPERIENCE — sensory observation (SEE / HEAR / FEEL)
# Slide 2: SAFETY — 3 situational A/B choices

ENV_OBSERVE={
    # see_q, hear_q, feel_q
    "森林":("树, 叶子, 花?",       "鸟? 风? 虫?",      "凉 还是 暖?"),
    "山地":("石头, 山, 天空?",     "风的声音?",         "冷 还是 热?"),
    "草地":("草, 花, 虫?",         "风? 蜜蜂? 鸟?",    "软 还是 硬? 暖?"),
    "河边":("水, 鱼, 石头?",       "水流的声音?",       "凉 还是 暖? 湿吗?"),
    "沙漠":("沙, 太阳, 仙人掌?",   "安静 还是 有风?",   "热 还是 冷? 干?"),
    "雪地":("雪, 冰, 天空?",       "踩雪的声音?",       "冷 还是 暖?"),
}

# Stage 2 (想 Think) — 4 visual hint chips per env to scaffold brainstorming
ENV_THINK={
    "森林":[("🍄","蘑菇?"),("🐍","蛇?"),("🐻","熊?"),("😱","迷路?")],
    "山地":[("🪨","滑倒?"),("⛈️","下雨?"),("🥶","太冷?"),("🌫️","雾?")],
    "草地":[("🐝","蜜蜂?"),("🌞","太阳晒?"),("🦟","小虫?"),("🌾","看不到脚?")],
    "河边":[("💦","滑倒?"),("🌊","掉水里?"),("🐍","水蛇?"),("🥶","水太凉?")],
    "沙漠":[("🥵","太热?"),("💧","没水?"),("🌪️","沙尘暴?"),("😵","迷路?")],
    "雪地":[("🥶","冻僵?"),("🧊","薄冰?"),("🌨️","暴风雪?"),("👁️","太亮?")],
}

# Stage 3 (说 Say) — teacher summary line
ENV_SAY={
    "森林":"森林里可能有: 不认识的蘑菇, 野生动物, 容易迷路。",
    "山地":"山上可能: 滑倒, 滚石, 突然下雨, 雷电。",
    "草地":"草地上有: 蜜蜂, 小虫, 太阳很大, 看不到脚下。",
    "河边":"河边: 地很滑, 水不能喝, 水深, 不能一个人下水。",
    "沙漠":"沙漠: 中午很热, 容易脱水, 容易迷路, 沙尘暴。",
    "雪地":"雪地: 太冷, 雪反光, 冰可能薄, 看不清路。",
}

# Stage 4 (判 Decide) — 3 A/B safety choices (correct: 0=A, 1=B)
ENV_DECIDE={
    "森林":[
        ("看到不认识的蘑菇", "摸一摸", "不要碰", 1),
        ("听到树丛里有声音", "跑开", "安静站着", 1),
        ("看到野生动物",     "走近看", "站住别动", 1),
    ],
    "山地":[
        ("山顶很滑",         "走快", "走慢", 1),
        ("看到石头掉下来",   "去摸", "跑开", 1),
        ("天上有黑云",       "继续上山", "下山", 1),
    ],
    "草地":[
        ("蜜蜂飞过来",       "挥手赶走", "不动", 1),
        ("太阳很大",         "戴帽子", "不戴", 0),
        ("草丛里有小虫",     "拍打", "不动", 1),
    ],
    "河边":[
        ("想喝水",           "河里的水", "水壶里的水", 1),
        ("湿石头要走过去",   "跳过去", "慢慢走", 1),
        ("想下水玩",         "自己下水", "等大人", 1),
    ],
    "沙漠":[
        ("中午太阳很大",     "继续走", "找阴凉", 1),
        ("渴了",             "喝水", "等等再喝", 0),
        ("走的时候",         "自己走", "跟着大人", 1),
    ],
    "雪地":[
        ("看到河面是冰",     "走过去", "绕道走", 1),
        ("太阳 + 雪很亮",    "戴墨镜", "不戴", 0),
        ("手很冷",           "戴手套", "不戴", 0),
    ],
}

# Stage 5 (做 Act) — movement prompt
ENV_ACT={
    "森林":"弯腰拨树枝走路 — 一步一步, 安静!",
    "山地":"扶着假岩壁, 一步一步慢慢爬。",
    "草地":"慢慢走 + 看脚下 + 戴帽子动作。",
    "河边":"假装小心踩湿石头, 双手张开。",
    "沙漠":"擦汗 + 慢慢走 + 喝水动作。",
    "雪地":"发抖 + 跺脚 + 抱紧自己。",
}

env_list=[
    ("🌲","森林","Forest",FOREST),
    ("🏔️","山地","Mountain",MOUNTAIN),
    ("🌾","草地","Grassland",GRASS),
    ("🏞️","河边","Riverside",RIVER),
    ("🏜️","沙漠","Desert",DESERT),
    ("❄️","雪地","Snow",SNOW),
]
for em,cn,en_name,env_color in env_list:
    see,hear,feel=ENV_OBSERVE[cn]
    s=experience_slide(em,cn,en_name,env_color,see,hear,feel,ENV_THINK[cn]);n+=1;pn(s,n)
    s=safety_slide(em,cn,en_name,env_color,ENV_SAY[cn],ENV_DECIDE[cn],ENV_ACT[cn]);n+=1;pn(s,n)

# ============================================================
# 12 COMPARISON TABLE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🗂️ 6 种环境对比  Compare 6 Environments",SUN)
tb(s,0.4,0.85,9,0.3,"每种环境都不一样，你能记住它们的危险吗？",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(4,7,Inches(0.3),Inches(1.2),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.4)
for i in range(1,7):t.columns[i].width=Inches(1.333)
rows=[
    ["","🌲 森林","🏔️ 山地","🌾 草地","🏞️ 河边","🏜️ 沙漠","❄️ 雪地"],
    ["✨ 特点","树多\n阴凉","陡高\n风大","平坦\n有花","有水\n有鱼","热/沙\n干燥","冷白\n有雪"],
    ["⚠️ 危险","迷路\n动物","滑倒\n天气","晒\n虫子","溺水","中暑\n迷路","冻伤\n薄冰"],
    ["✅ 规则","跟紧老师","穿好鞋","戴帽子","远离水","多喝水","穿厚衣"],
]
for r,rd in enumerate(rows):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True
        p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        rn=p.add_run();rn.text=ct.split('\n')[0];rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        for line in ct.split('\n')[1:]:
            p2=tf.add_paragraph();p2.alignment=PP_ALIGN.CENTER
            rn2=p2.add_run();rn2.text=line;rn2.font.name='KaiTi';rn2.font.size=Pt(11);rn2.font.color.rgb=DARK
        if r==0:
            rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=PINE
        elif c==0:
            rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        else:
            rn.font.color.rgb=DARK
            if r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF5)
pn(s,n)

# ============================================================
# 12b WHOLE-CLASS GAMES + Sentence frames (NEW)
# ============================================================
s=gesture_game_slide();n+=1;pn(s,n)
s=where_am_i_slide();n+=1;pn(s,n)

# ============================================================
# 13 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","复习 + 语言目标 (认字 + 写字)\n我会认 6 个词  ·  我会写 3 个词",SUN,"📖")
n+=1

# ============================================================
# 14 Quick review — danger match
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔄 快速复习  Quick Review — 找危险",SUN)
tb(s,0.4,0.85,9,0.3,"把环境和它最大的危险连起来 (口头)  Match the environment with its biggest danger!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
envs_q=[("🌲","森林",FOREST),("🏔️","山地",MOUNTAIN),("🌾","草地",GRASS),("🏞️","河边",RIVER),("🏜️","沙漠",DESERT),("❄️","雪地",SNOW)]
dangers_q=["🥵 中暑","🌊 溺水","🦟 虫咬","🪨 滑倒","🥶 冻伤","🐻 动物"]
# Left col: environments
for i,(em,cn,cl) in enumerate(envs_q):
    y=1.3+i*0.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(y),Inches(3.4),Inches(0.55))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.95,y+0.08,3.2,0.4,f"{em} {cn}",sz=16,b=True,c=WHITE)
# Right col: dangers (shuffled order)
for i,d in enumerate(dangers_q):
    y=1.3+i*0.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.8),Inches(y),Inches(3.4),Inches(0.55))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ALERT;sh.line.width=Pt(2)
    tb(s,5.95,y+0.08,3.2,0.4,d,sz=15,b=True,c=ALERT)
tb(s,4.25,2.9,1.5,0.4,"?",sz=40,b=True,c=SUN,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 15-20 我会认 word cards
# ============================================================
read_words=[
    ("森林","sēn lín","forest","森林里有很多大树。","📷 森林"),
    ("山地","shān dì","mountain","爬山要小心滑倒。","📷 山地"),
    ("草地","cǎo dì","grassland","草地上有很多花和虫。","📷 草地"),
    ("河边","hé biān","riverside","河边很危险，不能独自去。","📷 河边"),
    ("沙漠","shā mò","desert","沙漠里很热，要多喝水。","📷 沙漠"),
    ("雪地","xuě dì","snow","雪地里很冷，要穿厚衣服。","📷 雪地"),
]
for w,py,en,sent,img in read_words:
    s=word_card_read(w,py,en,sent,img);n+=1;pn(s,n)

# ============================================================
# 21 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",SUN)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","把字卡贴在\n白板上，老师\n说词语，学生拍！",WARM),
    ("2️⃣","举牌游戏\nShow Me","每人 6 张字卡\n老师说词语\n举正确的卡",RGBColor(0xFF,0xF3,0xE0)),
    ("3️⃣","抢椅子\nMusical Chairs","椅子上放字卡\n音乐停，读出词",RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣","传话筒\nPass the Mic","传球，停下的人\n读字卡并造句",RGBColor(0xE3,0xF2,0xFD)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 22-24 我会写 cards
# ============================================================
write_words=[
    ("森林","sēn lín","forest","📷 森林"),
    ("山地","shān dì","mountain","📷 山地"),
    ("河边","hé biān","riverside","📷 河边"),
]
for w,py,en,img in write_words:
    s=word_card_write(w,py,en,img);n+=1;pn(s,n)

# ============================================================
# 25 SESSION 3 DIVIDER
# ============================================================
div("Session 3  下午","写 Booklet + 动手做项目\n🖼️ 自然拼贴画  ·  🚧 安全标志  ·  🎭 我演你猜",BROWN,"🎒")
n+=1

# ============================================================
# 26 Booklet
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成"认识自然"练习册  Day 1 Booklet',BROWN)
ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图 / Booklet pages")
pn(s,n)

# ============================================================
# 27 Projects overview
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 动手时间！  Hands-On Time — 3 个活动",BROWN)
projects=[
    ("PROJECT 1","🖼️ 自然拼贴画","Nature Texture Collage","用树叶、石子、纸巾\n拼一幅自然环境",WARM,PINE),
    ("PROJECT 2","🚧 安全标志设计","Safety Sign Design","用彩笔和形状纸\n设计一个安全标志",RGBColor(0xFF,0xE0,0xB2),SUN),
    ("ACTIVITY","🎭 我演你猜","Charades","抽环境卡\n只用动作表演！",RGBColor(0xDC,0xED,0xC8),GREEN_OK),
]
for i,(lbl,nm,en,d,bgc,cl) in enumerate(projects):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(3.1),Inches(4.15))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,x+0.1,1.05,2.9,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.4,2.9,0.6,nm,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.0,2.9,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.2,2.45,2.8,1.6,"📷 示范")
    ls=d.split('\n')
    tf=tb(s,x+0.15,4.15,2.85,0.5,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for ln in ls[1:]:ap(tf,ln,sz=12,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 28 Project 1 — Nature Texture Collage
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🖼️ Project 1: 自然拼贴画  Nature Collage",PINE)
# Left: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=PINE;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.3,"🍂 干树叶  Dry leaves",sz=13,c=DARK)
ap(tf,"🪵 小树枝  Small twigs",sz=13,c=DARK)
ap(tf,"🪨 石子  Pebbles",sz=13,c=DARK)
ap(tf,"🧻 纸巾 (揉成雪)  Tissue (snow)",sz=13,c=DARK)
ap(tf,"🌊 铝箔纸 (代表河流)  Foil (river)",sz=13,c=DARK)
ap(tf,"💧 胶水  Glue",sz=13,c=DARK)
# Right: steps
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.3,"1️⃣ 选一个自然环境 (森林/河边/沙漠/雪地)",sz=13,c=DARK)
ap(tf2,"2️⃣ 用不同材料拼贴出这个环境",sz=13,c=DARK)
ap(tf2,"3️⃣ 贴牢、晾一下",sz=13,c=DARK)
ap(tf2,"4️⃣ 向同学介绍自己的作品",sz=13,c=DARK)
# Bottom: sentence frames
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.9),Inches(9.4),Inches(1.25))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=PINE;sh3.line.width=Pt(2)
tb(s,0.5,4.0,9,0.35,"🗣️ 展示句型  Say These:",sz=14,b=True,c=PINE)
tb(s,0.5,4.4,4.5,0.35,"· 这是森林。",sz=14,c=DARK)
tb(s,0.5,4.7,4.5,0.35,"· 这是河边。",sz=14,c=DARK)
tb(s,5.2,4.4,4.5,0.35,"· 这里有树叶、石头和小河。",sz=14,c=DARK)
pn(s,n)

# ============================================================
# 29 Project 2 — Safety Sign
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🚧 Project 2: 安全标志设计  Safety Sign",SUN)
# Left: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=PINE;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.5,"🖍️ 彩笔  Markers",sz=13,c=DARK)
ap(tf,"⭕ 圆形纸  Circle paper",sz=13,c=DARK)
ap(tf,"🔺 三角形纸  Triangle paper",sz=13,c=DARK)
ap(tf,"📎 胶水/胶带  Glue or tape",sz=13,c=DARK)
# Middle: example signs
sh_ex=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.1),Inches(4.4),Inches(0.4))
sh_ex.fill.solid();sh_ex.fill.fore_color.rgb=BROWN;sh_ex.line.fill.background()
tb(s,0.4,3.13,4.2,0.35,"💡 示例主题  Example Themes",sz=14,b=True,c=WHITE)
tf_ex=tb(s,0.4,3.6,4.4,1.5,"⚠️ 河边危险  Danger by river",sz=13,c=DARK)
ap(tf_ex,"🏃 不要乱跑  Don't run around",sz=13,c=DARK)
ap(tf_ex,"💧 小心滑倒  Slippery!",sz=13,c=DARK)
ap(tf_ex,"🛡️ 注意安全  Stay safe",sz=13,c=DARK)
# Right: steps + sentence frames
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=PINE;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.5,"1️⃣ 老师先示范几个安全标志",sz=13,c=DARK)
ap(tf2,"2️⃣ 选一个安全主题",sz=13,c=DARK)
ap(tf2,"3️⃣ 用圆形/三角形 + 彩笔设计",sz=13,c=DARK)
ap(tf2,"4️⃣ 说一说你的标志表示什么",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.1),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=GREEN_OK;sh3.line.fill.background()
tb(s,5.0,3.13,4.6,0.35,"🗣️ 展示句型  Say These",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.6,4.7,1.5,"· 这是我的安全标志。",sz=13,c=DARK)
ap(tf3,"· 它表示河边危险。",sz=13,c=DARK)
ap(tf3,"· 它告诉我们不要乱跑。",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 30 Activity — Charades rules
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 我演你猜  Charades — 规则 Rules",GREEN_OK)
rules_data=[
    ("1️⃣","上台抽卡","一名学生上台，抽一张环境卡"),
    ("2️⃣","只能表演","只能用动作表演，不能说话"),
    ("3️⃣","大家猜","其他同学根据动作猜环境"),
    ("4️⃣","一起说","猜对后，全班一起大声说！"),
]
for i,(num,t,d) in enumerate(rules_data):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.2),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GREEN_OK;sh.line.width=Pt(3)
    tb(s,x+0.1,1.1,2.0,0.6,num,sz=30,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.8,2.0,0.5,t,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.35,2.0,1.1,d,sz=12,c=DARK,a=PP_ALIGN.CENTER)
# Ask sentences
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.8),Inches(9.4),Inches(1.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=SUN;sh3.line.width=Pt(2)
tb(s,0.5,3.9,9,0.4,"🗣️ 可用提问句型  Ask:",sz=15,b=True,c=SUN)
tb(s,0.5,4.3,4.6,0.35,"· 这是哪里？",sz=14,c=DARK)
tb(s,0.5,4.65,4.6,0.35,"· 是森林吗？",sz=14,c=DARK)
tb(s,5.2,4.3,4.5,0.35,"· 我觉得是沙漠。",sz=14,c=DARK)
tb(s,5.2,4.65,4.5,0.35,"· 对不对？",sz=14,c=DARK)
pn(s,n)

# ============================================================
# 31 Charades hints
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 我演你猜  表演提示 Acting Hints",GREEN_OK)
hints=[
    ("🌲","森林","弯腰走路、拨开树枝",FOREST),
    ("🏔️","山地","小心走路、假装爬坡",MOUNTAIN),
    ("🌾","草地","轻松走路、看远方",GRASS),
    ("🏞️","河边","停下来、低头看水",RIVER),
    ("🏜️","沙漠","擦汗、感觉很热",DESERT),
    ("❄️","雪地","发抖、抱紧身体",SNOW),
]
for i,(em,cn,act,cl) in enumerate(hints):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=0.95+row*2.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.1,2.8,0.55,f"{em} {cn}",sz=19,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,y+0.7,2.7,1.1,act,sz=14,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 32 Day 1 Badge / Visa stamp
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.3,9,0.7,"🎖️ Day 1 探险家徽章  Explorer Badge",sz=24,b=True,c=PINE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.05),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PINE;sh.line.width=Pt(5)
tf=tb(s,3.6,1.30,2.8,2.7,"DAY 1",sz=18,b=True,c=SUN,a=PP_ALIGN.CENTER)
ap(tf,"🏕️",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"认识自然",sz=20,b=True,c=PINE,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=13,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
ap(tf,"🌲🏔️🌾🏞️🏜️❄️",sz=14,a=PP_ALIGN.CENTER)
# 6 stars filled in — earned!
sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.3),Inches(4.20),Inches(7.4),Inches(0.65))
sb.fill.solid();sb.fill.fore_color.rgb=SUNYEL;sb.line.color.rgb=SUN;sb.line.width=Pt(2.5)
tb(s,1.3,4.25,7.4,0.55,"⭐  ⭐  ⭐  ⭐  ⭐  ⭐",sz=32,b=True,c=ALERT,a=PP_ALIGN.CENTER)
tb(s,1,4.95,8,0.4,"6 颗星都拿到啦! All 6 stars earned! 🎉",sz=16,b=True,c=PINE,a=PP_ALIGN.CENTER)
tb(s,1,5.30,8,0.3,"学会了 6 种环境 · 6 条安全规则 · 3 个词",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 33 Tomorrow preview
# ============================================================
s=ns();n+=1;bg(s,PINE)
tb(s,0.5,0.9,9,0.8,"🔭 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.5,"Day 2 — 野外工具与装备",sz=28,b=True,c=SUN,a=PP_ALIGN.CENTER)
ap(tf,"Wilderness Tools & Gear",sz=16,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🎒 背包里要带什么？",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"What do explorers pack?",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，小探险家！",sz=15,c=WARM,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/野外生存与探险wilderness_pbl/day1_nature.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
