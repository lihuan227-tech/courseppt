#!/usr/bin/env python3
"""
小小艺术家 Little Artist Unit — Day 4: 当代艺术 大胆表达 Contemporary Art · Bold Expression
Bold contemporary palette: hot pink + Kusama dots + op-art black & white.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Bold Contemporary ---
HOT_PINK      = RGBColor(0xFF,0x1B,0x6B)   # primary — Kusama dot pink
DOT_RED       = RGBColor(0xE8,0x2A,0x2A)   # Kusama signature red
JET           = RGBColor(0x14,0x14,0x14)   # high-contrast black for op-art
SNOW          = RGBColor(0xFA,0xFA,0xFA)   # off-white for op-art
LIME          = RGBColor(0xCD,0xDC,0x39)   # pop yellow-green
ELECTRIC_BLUE = RGBColor(0x00,0xB0,0xFF)   # electric accent
# carry-overs from D1 for variety
MAGENTA = RGBColor(0xD8,0x1B,0x60)
YELLOW  = RGBColor(0xFF,0xC1,0x07)
CORAL   = RGBColor(0xFF,0x70,0x43)
PURPLE  = RGBColor(0x9C,0x27,0xB0)
GREEN_L = RGBColor(0x66,0xBB,0x6A)
CREAM   = RGBColor(0xFF,0xF8,0xE7)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xFF,0xF3,0xE0)
IMGBG   = RGBColor(0xF0,0xE8,0xF0)
GREEN_OK= RGBColor(0x38,0x8E,0x3C)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=HOT_PINK,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    palette=[YELLOW,LIME,ELECTRIC_BLUE,DOT_RED,HOT_PINK,SNOW]
    dot_colors=[c for c in palette if c!=color][:4]
    positions=[(0.8,4.7),(1.6,4.5),(7.8,4.5),(8.6,4.7)]
    for (x,y),cl in zip(positions,dot_colors):
        d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(0.4),Inches(0.4))
        d.fill.solid();d.fill.fore_color.rgb=cl;d.line.fill.background()
    return s

def dot(s,x,y,r,color):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(r),Inches(r))
    d.fill.solid();d.fill.fore_color.rgb=color;d.line.fill.background()

def example_slide(cat_emoji, cat_cn, work_cn, work_en, artist, fact, img_lb, color):
    s=ns();bg(s,CREAM)
    hb(s,f"{cat_emoji} {cat_cn}  ·  《{work_cn}》",color)
    ib(s,0.5,0.95,9,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,work_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,artist,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def example_slide_generic(cat_emoji, cat_cn, name_cn, name_en, sub_label, fact, img_lb, color):
    s=ns();bg(s,CREAM)
    hb(s,f"{cat_emoji} {cat_cn}  ·  {name_cn}",color)
    ib(s,0.5,0.95,9,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,name_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,sub_label,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def type_overview_slide(emoji, name_cn, name_en, color, hint, subtypes):
    s=ns();bg(s,CREAM);hb(s,f"{emoji} {name_cn}  {name_en}",color)
    tb(s,0.4,0.85,9,0.35,hint,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    cols = 3 if len(subtypes)<=6 else 4
    rows = (len(subtypes)+cols-1)//cols
    card_w = (9.4 - 0.15*(cols-1))/cols
    card_h = 1.7 if rows<=2 else 1.45
    y_start = 1.3 if rows<=2 else 1.25
    y_step = card_h + 0.25 if rows<=2 else card_h + 0.2
    for i,(sem,scn,sen) in enumerate(subtypes):
        col=i%cols;row=i//cols
        x=0.3+col*(card_w+0.15);y=y_start+row*y_step
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(card_w),Inches(card_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
        tb(s,x+0.1,y+0.1,card_w-0.2,0.55,sem,sz=30,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.75,card_w-0.2,0.4,scn,sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+1.15,card_w-0.2,0.3,sen,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    return s

def word_card_read(w,py,en,sent,img,color=HOT_PINK):
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=color,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=color)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,img,color=HOT_PINK):
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=color)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=color)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    return s

n=0

# ============================================================
# 1 COVER — Hot pink + Kusama-dot motif
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.7,"Contemporary Art · Bold!",sz=30,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.45,"当代艺术 · 大胆表达",sz=22,c=HOT_PINK,a=PP_ALIGN.CENTER)
# Big op-art badge: hot pink circle, jet ring, jet/lime dots scattered
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.55),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.color.rgb=JET;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.85),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=SNOW;sh2.line.color.rgb=JET;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.0,2.8,0.4,"DAY 4",sz=16,b=True,c=DOT_RED,a=PP_ALIGN.CENTER)
ap(tf,"🎨",sz=52,a=PP_ALIGN.CENTER)
ap(tf,"当代艺术",sz=20,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
ap(tf,"BOLD EXPRESSION",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Kusama-dot scatter — varied sizes around the badge
for x,y,r,c in [(1.0,1.6,0.55,DOT_RED),(1.6,3.5,0.4,JET),(0.8,4.5,0.5,LIME),
                (8.0,1.4,0.5,JET),(8.5,3.3,0.6,DOT_RED),(7.6,4.6,0.4,ELECTRIC_BLUE),
                (2.3,2.8,0.25,JET),(7.4,2.6,0.3,JET),(1.3,2.4,0.22,DOT_RED),(8.4,4.2,0.22,LIME)]:
    dot(s,x,y,r,c)
tb(s,1,5.05,8,0.4,"🟣 这是艺术吗？为什么？  Is this art? Why?",sz=14,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","什么是当代艺术？认识 5 种类型",HOT_PINK),
    ("Session 2  下午","2:00-2:45","复习 + 我会认 5 词 · 我会写 3 词",ELECTRIC_BLUE),
    ("Session 3  下午","3:00-4:30","Project：视错觉画 / 黑白对称剪纸",DOT_RED)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=WARM)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"🎨 内容目标  Content:",sz=20,b=True,c=HOT_PINK)
tf=tb(s,0.7,1.4,9,1.4,"1. 了解当代艺术的特点（不一定像真、可以大胆、可以用生活物品）",sz=14,c=DARK)
ap(tf,"2. 认识 5 种当代艺术：超现实 / 生活物品 / 草间弥生 / 视错觉 / 黑白对称",sz=14,c=DARK)
ap(tf,"3. 学习用点、线、对称、重复来创作",sz=14,c=DARK)
ap(tf,"4. 学会观察和讨论当代艺术",sz=14,c=DARK)
tb(s,0.5,3.05,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=ELECTRIC_BLUE)
tb(s,0.7,3.55,9,0.4,"👀 我会认：点  线  黑白  对称  花纹",sz=15,b=True,c=DARK)
tb(s,0.7,4.05,9,0.4,"✍️ 我会写：点  线  花纹",sz=15,b=True,c=DARK)
tb(s,0.5,4.65,9,0.5,"🖌️ 实践目标：完成 D4 booklet + 1 幅作品（视错觉 OR 对称剪纸）",sz=14,c=DOT_RED)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","什么是当代艺术？  What is Contemporary Art?\n🎨 奇怪 · 大胆 · 用任何东西都行！",HOT_PINK,"🟣")
n+=1

# ============================================================
# 5 GALLERY INTRO — preview the 6 examples we'll see one by one
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🖼️ 看一看  Look at These Artworks",HOT_PINK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.0),Inches(8.0),Inches(1.5))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,1.2,1.15,7.6,0.55,"接下来，我们一起看 6 幅画",sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.2,1.7,7.6,0.4,"Next, we'll look at 6 artworks together",sz=14,c=LIME,a=PP_ALIGN.CENTER)
tb(s,1.2,2.1,7.6,0.4,"它们都是「当代艺术」 — 你能发现什么共同点？",sz=14,c=YELLOW,a=PP_ALIGN.CENTER)
# 6 small thumbnails preview (2 rows × 3)
items=[("🕰️","融化的钟"),("🚭","这不是烟斗"),("🎃","圆点南瓜"),
       ("👁️","视错觉"),("🚲","自行车轮"),("🦋","黑白对称")]
for i,(em,t) in enumerate(items):
    col=i%3;row=i//3
    x=0.6+col*3.0;y=2.7+row*1.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.8),Inches(1.05))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=HOT_PINK;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.05,0.9,0.55,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x+1.0,y+0.1,1.7,0.4,t,sz=12,b=True,c=HOT_PINK)
    tb(s,x+1.0,y+0.5,1.7,0.4,f"#{i+1}",sz=10,c=GRAY)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.15),Inches(9.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=JET;sh.line.fill.background()
tb(s,0.4,5.18,9.2,0.35,"👀 仔细看 — 不要急着回答，让眼睛先告诉你！",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 5a-5f  ONE ARTWORK PER SLIDE (big image, no features revealed)
# ============================================================
gallery=[
    ("📷 达利《永恒的记忆》(融化的钟)  Dalí: The Persistence of Memory","《永恒的记忆》(融化的钟) · Salvador Dalí","🕰️",1),
    ("📷 马格利特《这不是烟斗》  Magritte: The Treachery of Images","《这不是烟斗》· René Magritte","🚭",2),
    ("📷 草间弥生 圆点南瓜  Yayoi Kusama: Pumpkin (黄底黑点)","圆点南瓜 · 草间弥生 Yayoi Kusama","🎃",3),
    ("📷 Bridget Riley 黑白条纹 (视错觉)  Op-Art lines","视错觉条纹 · Bridget Riley","👁️",4),
    ("📷 杜尚 自行车轮  Duchamp: Bicycle Wheel (现成品)","现成品 — 自行车轮 · Marcel Duchamp","🚲",5),
    ("📷 黑白对称图案  Symmetric paper-cut pattern","黑白对称图案 · Symmetry","🦋",6),
]
for img_lb,title,em,idx in gallery:
    s=ns();n+=1;bg(s,CREAM);hb(s,f"🖼️ 看 — 第 {idx} 幅 / {len(gallery)}  Look · {idx}/{len(gallery)}",HOT_PINK)
    ib(s,0.5,1.0,9.0,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.6),Inches(9.0),Inches(0.45))
    sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
    tb(s,0.6,4.65,1.0,0.35,em,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.6,4.68,7.8,0.35,title,sz=14,b=True,c=WHITE)
    tb(s,0.5,5.15,9.0,0.35,"👀 像不像真的？奇怪吗？— 你看到了什么？",sz=12,c=HOT_PINK,a=PP_ALIGN.CENTER)
    pn(s,n)

# ============================================================
# 6 INQUIRY — what do you see? (Student observation, no answers)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🤔 你看到什么？  What Did You See?",HOT_PINK)
tb(s,0.4,0.9,9.2,0.35,"看完 6 张画，一起想一想 / After looking, let's think together",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# LEFT: 5 observation questions
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.35),Inches(5.3),Inches(3.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=HOT_PINK;sh.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.35),Inches(5.3),Inches(0.5))
head.fill.solid();head.fill.fore_color.rgb=HOT_PINK;head.line.fill.background()
tb(s,0.45,1.43,5.0,0.4,"❓ 一起想一想  Let's Think",sz=14,b=True,c=WHITE)
qs=[
    "画得像不像真实的东西？",
    "看起来奇怪 还是 普通？",
    "用了点点 / 线条 / 重复图案吗？",
    "有没有「日常物品」也变成了艺术？",
    "你觉得这是「艺术」吗？为什么？",
]
tf=tb(s,0.5,2.0,5.0,0.5,f"❓ {qs[0]}",sz=13,c=DARK)
for q in qs[1:]:
    ap(tf,"",sz=6)
    ap(tf,f"❓ {q}",sz=13,c=DARK)
# RIGHT: sentence frames
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.85),Inches(1.35),Inches(3.85),Inches(3.7))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=ELECTRIC_BLUE;sh.line.width=Pt(2.5)
head2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.85),Inches(1.35),Inches(3.85),Inches(0.5))
head2.fill.solid();head2.fill.fore_color.rgb=ELECTRIC_BLUE;head2.line.fill.background()
tb(s,6.0,1.43,3.55,0.4,"💬 我会说  I Can Say",sz=14,b=True,c=WHITE)
frames=[
    "我看到 ________。",
    "它和真实 ________。",
    "我觉得很 ________。",
    "用了 ________。",
    "我喜欢 ________ 因为 ________。",
]
tf2=tb(s,6.05,2.0,3.65,0.5,f"·  {frames[0]}",sz=12,c=DARK)
for fr in frames[1:]:
    ap(tf2,"",sz=6)
    ap(tf2,f"·  {fr}",sz=12,c=DARK)
# Bottom transition
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.15),Inches(9.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=JET;sh.line.fill.background()
tb(s,0.4,5.18,9.2,0.35,"👉 大家说完，我们一起总结 — 这就是「当代艺术」!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 7 REVEAL — 你说对啦！big idea (now confirms what students discovered)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💡 你说对啦！当代艺术 = ?",HOT_PINK)
tb(s,0.4,0.9,9.2,0.35,"You discovered it! Here's the big idea.",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.2),Inches(1.4),Inches(7.6),Inches(1.5))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,1.4,1.5,7.2,0.7,"艺术不一定要像真的！",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.4,2.2,7.2,0.55,"Art doesn't have to look REAL!",sz=18,c=LIME,a=PP_ALIGN.CENTER)
bubbles=[
    ("🤪","可以很奇怪","Can be WEIRD",DOT_RED),
    ("🌈","可以很大胆","Can be BOLD",ELECTRIC_BLUE),
    ("🙋","每个人都可以做","Anyone can do it",LIME),
]
for i,(em,cn,en,cl) in enumerate(bubbles):
    x=0.5+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.2),Inches(3.0),Inches(1.95))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,3.3,2.8,0.65,em,sz=40,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.0,2.8,0.45,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.55,2.8,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 6 OLD vs NEW comparison
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🆚 古典 vs 当代  Classical vs Contemporary",HOT_PINK)
# Left card — classical
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(4.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PURPLE;sh.line.width=Pt(3)
tb(s,0.5,1.1,4.3,0.5,"🖼️ 古典艺术",sz=22,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.5,1.6,4.3,0.35,"Classical Art",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.6,2.05,4.1,1.5,"📷 蒙娜丽莎 / 古典油画")
tf=tb(s,0.6,3.65,4.1,0.4,"· 像照片，很真实",sz=14,c=DARK)
ap(tf,"· 画名画、画人、画风景",sz=14,c=DARK)
ap(tf,"· 几百年前的艺术",sz=14,c=DARK)
# Right card — contemporary
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.0),Inches(4.5),Inches(4.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=HOT_PINK;sh.line.width=Pt(3)
tb(s,5.2,1.1,4.3,0.5,"🟣 当代艺术",sz=22,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
tb(s,5.2,1.6,4.3,0.35,"Contemporary Art",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,5.3,2.05,4.1,1.5,"📷 草间弥生圆点 / 视错觉")
tf=tb(s,5.3,3.65,4.1,0.4,"· 奇怪、抽象、大胆",sz=14,c=DARK)
ap(tf,"· 可以用任何东西",sz=14,c=DARK)
ap(tf,"· 现在的艺术家也在做！",sz=14,c=DARK)
pn(s,n)

# ============================================================
# 7 FEATURES GRID 2x3
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📌 总结 — 当代艺术的 6 个特点  6 Features",HOT_PINK)
tb(s,0.4,0.9,9,0.4,"我们一起发现的 6 个特点 / The 6 features we found together",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
features=[
    ("🚫","不一定像真","Not realistic",HOT_PINK),
    ("🤪","可以奇怪","Can be weird",DOT_RED),
    ("💥","可以大胆","Can be bold",ELECTRIC_BLUE),
    ("🔁","点线形状重复","Repeat patterns",LIME),
    ("🪑","生活物品也是艺术","Everyday = art",PURPLE),
    ("🙋","自己的方式表达","Your own way",CORAL),
]
for i,(em,cn,en,cl) in enumerate(features):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.4+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.1,2.8,0.7,em,sz=38,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.95,2.8,0.45,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.4,2.8,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8 OVERVIEW — 5 types of contemporary art
# ============================================================
s=type_overview_slide("🎨","5 种当代艺术","5 Types of Contemporary Art",HOT_PINK,
    "今天我们来认识 5 种 — 你最喜欢哪一种？",
    [("🐟","超现实","Surrealism"),
     ("🚽","生活物品","Readymades"),
     ("🔴","草间弥生","Yayoi Kusama"),
     ("👁️","视错觉","Optical Illusion"),
     ("🦋","黑白对称","Symmetry")]);n+=1;pn(s,n)

# ============================================================
# 9 超现实 inquiry
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🐟 超现实  Surrealism — 想象的世界",PURPLE)
ib(s,0.4,0.95,4.6,3.4,"📷 鱼在天上飞 / 房子长在树上")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(0.95),Inches(4.5),Inches(3.4))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PURPLE;sh.line.width=Pt(2.5)
tb(s,5.4,1.1,4.1,0.5,"什么是超现实？",sz=20,b=True,c=PURPLE)
tf=tb(s,5.4,1.65,4.1,0.4,"· 把想象变成画",sz=14,c=DARK)
ap(tf,"· 不可能的事 → 画出来",sz=14,c=DARK)
ap(tf,"· 鱼在天上飞、房子长在树上",sz=14,c=DARK)
ap(tf,"· 像做梦一样的世界",sz=14,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.5),Inches(9.3),Inches(0.85))
sh2.fill.solid();sh2.fill.fore_color.rgb=PURPLE;sh2.line.fill.background()
tb(s,0.6,4.55,9.0,0.4,"🤔 如果你画一幅奇怪的画，你想画什么？",sz=18,b=True,c=WHITE)
tb(s,0.6,4.95,9.0,0.35,"If you painted a weird picture, what would you paint?",sz=12,c=WARM)
pn(s,n)

# ============================================================
# 10 超现实 example — Magritte 这不是烟斗
# ============================================================
s=example_slide("🐟","超现实","这不是烟斗","The Treachery of Images",
    "马格利特 René Magritte · 比利时 · 1929",
    "明明是烟斗，他说不是 — 你觉得呢？\nIt looks like a pipe… but is it?",
    "📷 马格利特 这不是烟斗",PURPLE);n+=1;pn(s,n)

# ============================================================
# 11 超现实 example — Dalí 永恒的记忆
# ============================================================
s=example_slide("🐟","超现实","永恒的记忆","The Persistence of Memory",
    "达利 Salvador Dalí · 西班牙 · 1931",
    "钟会化掉？时间在融化。\nClocks melting — time is melting too!",
    "📷 达利 融化的钟",PURPLE);n+=1;pn(s,n)

# ============================================================
# 12 生活物品 inquiry — Duchamp
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🚽 生活物品  Readymades — 这也是艺术？",ELECTRIC_BLUE)
ib(s,0.4,0.95,4.6,3.4,"📷 杜尚《泉》(陶瓷)")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(0.95),Inches(4.5),Inches(3.4))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ELECTRIC_BLUE;sh.line.width=Pt(2.5)
tb(s,5.4,1.1,4.1,0.5,"杜尚《泉》Fountain",sz=20,b=True,c=ELECTRIC_BLUE)
tb(s,5.4,1.55,4.1,0.35,"Marcel Duchamp · 1917",sz=12,c=GRAY)
tf=tb(s,5.4,1.95,4.1,0.4,"· 他把一个陶瓷物品放进美术馆",sz=14,c=DARK)
ap(tf,"· 没有画、没有雕，只是「放」",sz=14,c=DARK)
ap(tf,"· 大家都吵：这是艺术吗？",sz=14,c=DARK)
ap(tf,"· 100 年后 — 还在讨论！",sz=14,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.5),Inches(9.3),Inches(0.85))
sh2.fill.solid();sh2.fill.fore_color.rgb=ELECTRIC_BLUE;sh2.line.fill.background()
tb(s,0.6,4.55,9.0,0.4,"🤔 你身边有什么 — 也可以是艺术吗？",sz=18,b=True,c=WHITE)
tb(s,0.6,4.95,9.0,0.35,"Look around — what could be art?",sz=12,c=WARM)
pn(s,n)

# ============================================================
# 13 生活物品 — class brainstorm
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧺 来找找！  Look Around the Room",ELECTRIC_BLUE)
tb(s,0.4,0.85,9,0.4,"把它们摆在一起 = 艺术？  Arrange them together = art?",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
items=[
    ("🪑","椅子","Chair",HOT_PINK),
    ("🥤","杯子","Cup",ELECTRIC_BLUE),
    ("👟","鞋子","Shoe",DOT_RED),
    ("🧸","玩具","Toy",LIME),
    ("🍃","树叶","Leaves",GREEN_L),
    ("📰","报纸","Newspaper",PURPLE),
]
for i,(em,cn,en,cl) in enumerate(items):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.35+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.1,2.8,0.7,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.9,2.8,0.45,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.35,2.8,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 14 草间弥生 intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔴 草间弥生  Yayoi Kusama — 圆点女王",DOT_RED)
ib(s,0.4,0.95,4.6,3.8,"📷 草间弥生肖像 (戴红色波点假发)")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(0.95),Inches(4.5),Inches(3.8))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=DOT_RED;sh.line.width=Pt(3)
tb(s,5.4,1.1,4.1,0.5,"草间弥生 Yayoi Kusama",sz=20,b=True,c=DOT_RED)
tb(s,5.4,1.55,4.1,0.35,"日本艺术家 Japan · 1929 - 至今",sz=12,c=GRAY)
tf=tb(s,5.4,2.0,4.1,0.4,"· 用「圆点」表达自己",sz=14,c=DARK)
ap(tf,"· 从小就「看到」圆点",sz=14,c=DARK)
ap(tf,"· 圆点让她安心、不害怕",sz=14,c=DARK)
ap(tf,"· 现在 90 多岁，还在画！",sz=14,c=DARK)
ap(tf,"· 全世界最有名的奶奶艺术家",sz=14,c=DOT_RED,b=True)
# decorative dots
for x,y,r,c in [(0.6,5.05,0.18,DOT_RED),(1.0,5.15,0.12,JET),(1.4,5.0,0.15,DOT_RED),(2.0,5.1,0.1,JET),(2.4,5.05,0.16,DOT_RED)]:
    dot(s,x,y,r,c)
pn(s,n)

# ============================================================
# 15 草间弥生 example — 南瓜雕塑
# ============================================================
s=example_slide("🎃","草间弥生","南瓜","Pumpkin",
    "草间弥生 Yayoi Kusama · 1994 · 日本直岛",
    "黄底黑点的南瓜 · 她从小就看到圆点 · 圆点让她安心。\nYellow + black dots — dots make her feel safe.",
    "📷 草间弥生 南瓜雕塑 (黄+黑点)",DOT_RED);n+=1;pn(s,n)

# ============================================================
# 16 草间弥生 — Infinity Room
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔴 草间弥生  ·  无限镜屋  Infinity Mirror Room",DOT_RED)
ib(s,0.5,0.95,9,3.4,"📷 Infinity Room — 圆点 + 镜子 + 灯光")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.45),Inches(9),Inches(0.95))
sh.fill.solid();sh.fill.fore_color.rgb=DOT_RED;sh.line.fill.background()
tb(s,0.7,4.5,9.0,0.4,"走进去像在宇宙里 · 圆点无穷无尽",sz=18,b=True,c=WHITE)
tb(s,0.7,4.9,9.0,0.4,"🤔 如果你的房间全是圆点，你会怎样？  How would you feel?",sz=13,c=WARM)
pn(s,n)

# ============================================================
# 17 草间弥生 try-it inquiry
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 试一试  Try It! — 你来画圆点",DOT_RED)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.2),Inches(1.0),Inches(7.6),Inches(1.4))
sh.fill.solid();sh.fill.fore_color.rgb=DOT_RED;sh.line.fill.background()
tb(s,1.4,1.1,7.2,0.6,"选一个颜色 + 重复画圆点",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.4,1.75,7.2,0.5,"Pick ONE color · Repeat dots everywhere!",sz=16,c=WARM,a=PP_ALIGN.CENTER)
# 4 idea cards
ideas=[
    ("📄","纸上","On paper",DOT_RED),
    ("👕","衣服上","On clothes",ELECTRIC_BLUE),
    ("🪑","椅子上","On chair",LIME),
    ("✋","手上","On your hand",HOT_PINK),
]
for i,(em,cn,en,cl) in enumerate(ideas):
    x=0.5+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.7),Inches(2.1),Inches(1.8))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,2.8,1.9,0.7,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.55,1.9,0.45,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.05,1.9,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.75,9,0.4,"🤔 你想画在哪里？  Where do YOU want to dot?",sz=15,b=True,c=DOT_RED,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 18 视错觉 Optical Illusion intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👁️ 视错觉  Optical Illusion — 看起来在动！",JET)
ib(s,0.4,0.95,4.6,3.8,"📷 黑白同心圆 / 平行线视错觉")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(0.95),Inches(4.5),Inches(3.8))
sh.fill.solid();sh.fill.fore_color.rgb=SNOW;sh.line.color.rgb=JET;sh.line.width=Pt(3)
tb(s,5.4,1.1,4.1,0.5,"什么是视错觉？",sz=22,b=True,c=JET)
tb(s,5.4,1.65,4.1,0.35,"Optical Illusion",sz=12,c=GRAY)
tf=tb(s,5.4,2.05,4.1,0.4,"· 看起来在动",sz=15,c=JET)
ap(tf,"· 但其实没动！",sz=15,c=JET)
ap(tf,"· 是眼睛被「骗」了",sz=15,c=JET)
ap(tf,"· 黑白线条 + 重复 = 晃眼",sz=15,c=JET)
ap(tf,"",sz=8)
ap(tf,"🤔 盯着看 — 它在动吗？",sz=14,b=True,c=DOT_RED)
pn(s,n)

# ============================================================
# 19 视错觉 example — Bridget Riley
# ============================================================
s=example_slide_generic("👁️","视错觉","Bridget Riley · 黑白线","Black & White Lines",
    "英国艺术家 British · 1931 - 至今",
    "黑白线条让眼睛「晃」 — 盯着看，它在动吗？\nLook hard — does it move?",
    "📷 Bridget Riley 黑白波纹",JET);n+=1;pn(s,n)

# ============================================================
# 20 视错觉 example — Vasarely
# ============================================================
s=example_slide_generic("👁️","视错觉","Vasarely · 彩色方块","Colored Squares",
    "维克托·瓦萨雷利 Hungary/France · 1906-1997",
    "彩色方块 · 看起来鼓起来 — 中间是凹下去还是凸起来？\nDoes the center push out or sink in?",
    "📷 Vasarely 彩色立体方块",ELECTRIC_BLUE);n+=1;pn(s,n)

# ============================================================
# 21 黑白对称 intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🦋 黑白对称  Symmetry — 对折 → 剪 → 打开",JET)
ib(s,0.4,0.95,4.6,3.8,"📷 对称剪纸示范 (折→剪→打开)")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(0.95),Inches(4.5),Inches(3.8))
sh.fill.solid();sh.fill.fore_color.rgb=SNOW;sh.line.color.rgb=JET;sh.line.width=Pt(3)
tb(s,5.4,1.1,4.1,0.5,"什么是对称？",sz=22,b=True,c=JET)
tb(s,5.4,1.65,4.1,0.35,"What is symmetry?",sz=12,c=GRAY)
tf=tb(s,5.4,2.05,4.1,0.4,"· 左边 = 右边",sz=15,c=JET)
ap(tf,"· 像照镜子一样",sz=15,c=JET)
ap(tf,"· 对折 → 剪 → 打开 = 对称图案！",sz=15,c=JET)
ap(tf,"",sz=8)
ap(tf,"🤔 对称是什么？",sz=14,b=True,c=DOT_RED)
ap(tf,"   身体哪里对称？",sz=14,b=True,c=DOT_RED)
ap(tf,"   (眼睛、耳朵、手……)",sz=12,c=GRAY)
pn(s,n)

# ============================================================
# 22 黑白对称 examples — 4 thumbnails
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🦋 对称的例子  Symmetry Examples",JET)
tb(s,0.4,0.85,9,0.4,"看一看 — 这些都是对称的！",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
syms=[
    ("🦋","蝴蝶","Butterfly",HOT_PINK),
    ("❄️","雪花","Snowflake",ELECTRIC_BLUE),
    ("🌸","万花筒","Kaleidoscope",DOT_RED),
    ("✂️","剪纸","Paper Cut",LIME),
]
for i,(em,cn,en,cl) in enumerate(syms):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(2.2),Inches(3.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.45,2.0,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    ib(s,x+0.15,2.25,1.9,1.7,f"📷 {cn}图")
    tb(s,x+0.1,4.0,2.0,0.45,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.5,2.0,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 23 观察 + 讨论 — 4 mystery thumbnails
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔍 猜一猜  Guess Which is Which",HOT_PINK)
tb(s,0.4,0.85,9,0.4,"哪一幅是哪一种？说一说！  Which is which? Tell me!",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
mystery=[
    ("A","📷 鱼在天上飞","超现实？"),
    ("B","📷 圆点南瓜","草间弥生？"),
    ("C","📷 黑白波纹","视错觉？"),
    ("D","📷 蝴蝶图案","黑白对称？"),
]
for i,(lbl,img,q) in enumerate(mystery):
    col=i%2;row=i//2
    x=0.3+col*4.75;y=1.3+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.6),Inches(1.8))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=HOT_PINK;sh.line.width=Pt(2.5)
    tb(s,x+0.15,y+0.1,0.8,0.9,lbl,sz=36,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
    ib(s,x+1.05,y+0.15,2.1,1.5,img)
    tb(s,x+3.2,y+0.5,1.3,0.45,q,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+3.2,y+1.0,1.3,0.4,"举手猜！",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 24 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","复习 + 语言目标\n👀 我会认 5 个词  ·  ✍️ 我会写 3 个词",ELECTRIC_BLUE,"📖")
n+=1

# ============================================================
# 25 Quick review — match 5 features ↔ artwork
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔄 快速复习  Quick Review — 配一配",ELECTRIC_BLUE)
tb(s,0.4,0.85,9,0.35,"把词和图案连起来 (口头)",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
words=[("点",DOT_RED),("线",JET),("黑白",JET),("对称",HOT_PINK),("花纹",LIME)]
imgs=["📷 蝴蝶","📷 圆点南瓜","📷 黑白波纹","📷 直线条","📷 重复花纹"]
for i,(w,cl) in enumerate(words):
    y=1.2+i*0.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(y),Inches(3.4),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    txt_c = JET if cl==LIME else WHITE
    tb(s,0.75,y+0.1,3.2,0.5,w,sz=24,b=True,c=txt_c,a=PP_ALIGN.CENTER)
for i,im in enumerate(imgs):
    y=1.2+i*0.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.0),Inches(y),Inches(3.4),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ELECTRIC_BLUE;sh.line.width=Pt(2)
    tb(s,6.15,y+0.15,3.2,0.45,im,sz=14,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,4.3,3.0,1.4,0.5,"?",sz=44,b=True,c=DOT_RED,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 26-30 我会认 word cards
# ============================================================
read_words=[
    ("点","diǎn","dot","草间弥生喜欢画点。","📷 圆点南瓜",DOT_RED),
    ("线","xiàn","line","我用线条画图。","📷 黑白线条",JET),
    ("黑白","hēi bái","black & white","黑白图案很酷。","📷 黑白图案",JET),
    ("对称","duì chèn","symmetry","蝴蝶是对称的。","📷 蝴蝶",HOT_PINK),
    ("花纹","huā wén","pattern","我喜欢花纹的衣服。","📷 重复花纹",LIME),
]
for w,py,en,sent,img,cl in read_words:
    s=word_card_read(w,py,en,sent,img,cl);n+=1;pn(s,n)

# ============================================================
# 31 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",ELECTRIC_BLUE)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","老师说词\n学生拍正确的字卡",RGBColor(0xFF,0xF3,0xE0)),
    ("2️⃣","配图案\nPattern Match","看图案\n读对应的词",RGBColor(0xE3,0xF2,0xFD)),
    ("3️⃣","词语接龙\nWord Chain","一个接一个\n说花纹词",RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣","我画你猜\nDraw & Guess","画一画\n大家猜词",RGBColor(0xFC,0xE4,0xEC)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 32-34 我会写 word cards
# ============================================================
write_words=[
    ("点","diǎn","dot","📷 圆点",DOT_RED),
    ("线","xiàn","line","📷 线条",JET),
    ("花纹","huā wén","pattern","📷 重复花纹",LIME),
]
for w,py,en,img,cl in write_words:
    s=word_card_write(w,py,en,img,cl);n+=1;pn(s,n)

# ============================================================
# 35 SESSION 3 DIVIDER
# ============================================================
div("Session 3  下午","做大胆的艺术！  Make Bold Art!\n👁️ 视错觉画  ·  🦋 黑白对称剪纸",DOT_RED,"🎨")
n+=1

# ============================================================
# 36 Booklet
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成 D4 练习册  Day 4 Booklet',DOT_RED)
ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图 / Booklet pages")
pn(s,n)

# ============================================================
# 37 Projects overview
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 今天做什么？  Two Project Choices",DOT_RED)
tb(s,0.4,0.9,9,0.4,"老师选 OR 学生选！Teacher picks or you pick.",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
projects=[
    ("PROJECT 1","👁️ 视错觉画","Optical Illusion","用线 + 图案让画「看起来在动」\nLines + patterns make it move!",JET,JET),
    ("PROJECT 2","🦋 黑白对称剪纸","Symmetry Paper Cut","折 → 剪 → 打开 (Norton 风格)\nFold → cut → open!",HOT_PINK,HOT_PINK),
]
for i,(lvl,nm,en,d,ltcl,cl) in enumerate(projects):
    x=0.5+i*4.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(4.3),Inches(3.45))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(4)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.25),Inches(1.55),Inches(1.7),Inches(0.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=ltcl;sh2.line.fill.background()
    tb(s,x+0.3,1.6,1.6,0.4,lvl,sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,2.2,4.0,0.55,nm,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,2.8,4.0,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.5,3.15,3.4,1.1,"📷 作品示范")
    ls=d.split('\n')
    tf=tb(s,x+0.2,4.3,4.0,0.25,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Extension hint
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.95),Inches(9),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=LIME;sh.line.fill.background()
tb(s,0.7,5.0,8.6,0.4,"🌟 创作延伸：自己设计「重复图案」!  Extension: design your own pattern!",sz=14,b=True,c=JET,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 38 Project 1 — 视错觉画
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👁️ Project 1: 视错觉画  Optical Illusion",JET)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=JET;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"📄 黑白纸  Black & white paper",sz=13,c=DARK)
ap(tf,"🖊️ 黑笔 / 马克笔  Black pen / marker",sz=13,c=DARK)
ap(tf,"📐 直尺  Ruler",sz=13,c=DARK)
ap(tf,"⭕ 圆形纸  Round paper / template",sz=13,c=DARK)
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=LIME;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"💡 小贴士  Tips",sz=14,b=True,c=JET)
tf_k=tb(s,0.4,3.85,4.4,1.4,"· 线越细越平行，越像在动",sz=13,c=DARK)
ap(tf_k,"· 黑白对比最强烈",sz=13,c=DARK)
ap(tf_k,"· 可以加一点颜色 (一种就够)",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=DOT_RED;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 画一个中心 (一个点)",sz=13,c=DARK)
ap(tf2,"2️⃣ 从中心画放射线 / 同心圆",sz=13,c=DARK)
ap(tf2,"3️⃣ 重复 — 让它「晃」眼",sz=13,c=DARK)
ap(tf2,"4️⃣ 可以加一点颜色",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=ELECTRIC_BLUE;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"🗣️ 展示句型  Say These",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 这是我的视错觉画。",sz=13,c=DARK)
ap(tf3,"· 你看 — 它在动吗？",sz=13,c=DARK)
ap(tf3,"· 我用了点和线。",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 39 Project 2 — 对称剪纸
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🦋 Project 2: 黑白对称剪纸  Symmetry Paper Cut",HOT_PINK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"🟥 彩纸  Color paper",sz=13,c=DARK)
ap(tf,"✂️ 剪刀  Scissors",sz=13,c=DARK)
ap(tf,"🩹 胶水  Glue",sz=13,c=DARK)
ap(tf,"📄 白色衬纸  White backing paper",sz=13,c=DARK)
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=LIME;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"💡 小贴士  Tips",sz=14,b=True,c=JET)
tf_k=tb(s,0.4,3.85,4.4,1.4,"· 折 1 次 → 2 边对称",sz=13,c=DARK)
ap(tf_k,"· 折 2 次 → 4 边对称",sz=13,c=DARK)
ap(tf_k,"· 折 3 次 → 8 边对称 (像雪花！)",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=DOT_RED;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 把纸对折",sz=13,c=DARK)
ap(tf2,"2️⃣ 在折边画一半的图案",sz=13,c=DARK)
ap(tf2,"3️⃣ 沿线小心剪开",sz=13,c=DARK)
ap(tf2,"4️⃣ 打开 — 看对称！",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=ELECTRIC_BLUE;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"🗣️ 展示句型  Say These",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 这是我的剪纸。",sz=13,c=DARK)
ap(tf3,"· 你看，它是对称的！",sz=13,c=DARK)
ap(tf3,"· 左边 = 右边。",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 40 创作延伸 — 重复图案设计
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌟 创作延伸  Extension — 重复图案设计",LIME)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(9.2),Inches(1.2))
sh.fill.solid();sh.fill.fore_color.rgb=LIME;sh.line.fill.background()
tb(s,0.6,1.1,8.8,0.5,"用点 + 线 + 形状组合 → 画一个「重复花纹」",sz=24,b=True,c=JET,a=PP_ALIGN.CENTER)
tb(s,0.6,1.65,8.8,0.4,"Combine dots + lines + shapes → make a repeating pattern!",sz=14,c=JET,a=PP_ALIGN.CENTER)
# 3 sentence frame cards
frames=[
    ("🗣️","我用了 ____。","I used ___",HOT_PINK),
    ("🎨","这是我的花纹。","This is my pattern.",ELECTRIC_BLUE),
    ("🔁","你看，它在重复。","Look — it repeats!",DOT_RED),
]
for i,(em,cn,en,cl) in enumerate(frames):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.45),Inches(3.0),Inches(2.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.55,2.8,0.7,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.35,2.8,0.55,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.0,2.8,0.4,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.5,2.8,0.4,"☆ 跟我说",sz=11,c=cl,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 41 Gallery walk
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🖼️ 当代艺术展  Gallery Walk",HOT_PINK)
tb(s,0.4,0.9,9,0.5,"把作品贴在墙上，大家一起看！",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Hang your bold art and walk around like a museum!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("1️⃣","🖼️","贴作品\nHang it up"),
    ("2️⃣","👀","看同学的\nLook around"),
    ("3️⃣","🗣️","说一说\nTalk about it"),
    ("4️⃣","⭐","贴星星\nGive stars"),
]
colors_g=[HOT_PINK,ELECTRIC_BLUE,DOT_RED,LIME]
for i,(num,em,d) in enumerate(steps):
    x=0.3+i*2.4
    cl=colors_g[i]
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.1,2.0,0.5,num,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.6,2.0,0.8,em,sz=40,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,3.6,2.0,0.4,ls[0],sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 42 Day 4 Badge — Kusama-dot motif
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.4,9,0.8,"🎖️ Day 4 当代艺术家徽章  Contemporary Artist Badge",sz=24,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.color.rgb=JET;sh.line.width=Pt(5)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.75),Inches(1.65),Inches(2.5),Inches(2.5))
sh2.fill.solid();sh2.fill.fore_color.rgb=SNOW;sh2.line.color.rgb=JET;sh2.line.width=Pt(2)
tf=tb(s,3.85,1.85,2.3,2.5,"DAY 4",sz=16,b=True,c=DOT_RED,a=PP_ALIGN.CENTER)
ap(tf,"🎨",sz=36,a=PP_ALIGN.CENTER)
ap(tf,"当代艺术",sz=18,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=12,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
# Kusama-dot scatter
for x,y,r,c in [(2.4,2.0,0.32,DOT_RED),(7.1,2.0,0.32,JET),(2.2,3.5,0.28,JET),(7.3,3.5,0.28,DOT_RED),
                (1.7,2.8,0.22,DOT_RED),(7.8,2.8,0.22,DOT_RED),(2.8,1.5,0.18,JET),(6.9,1.5,0.18,JET),
                (2.6,4.2,0.2,LIME),(7.0,4.2,0.2,LIME)]:
    dot(s,x,y,r,c)
tb(s,1,4.55,8,0.4,"恭喜你完成 Day 4！You are bold! 🎉",sz=16,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"5 种当代艺术 · 5 个新词 · 1 幅大胆作品",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 43 Tomorrow preview
# ============================================================
s=ns();n+=1;bg(s,HOT_PINK)
tb(s,0.5,0.9,9,0.8,"🎨 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.6,"Day 5 — 我手画我心 · 艺术展览日",sz=24,b=True,c=LIME,a=PP_ALIGN.CENTER)
ap(tf,"Exhibition Day!",sz=16,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🖼️ 我们的大作品要诞生了！",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Our big artworks are coming!",sz=14,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，小艺术家！",sz=15,b=True,c=LIME,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/小小艺术家little_artist_pbl/day4_contemporary.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
