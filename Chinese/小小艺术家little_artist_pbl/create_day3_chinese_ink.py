#!/usr/bin/env python3
"""
小小艺术家 Little Artist Unit — Day 3: 中国水墨画 Chinese Ink Painting
Soft, traditional ink-painting palette: deep ink, jade, vermillion seal red, aged rice paper.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Traditional Ink Painting ---
INK        = RGBColor(0x1F,0x1F,0x1F)  # 浓墨 deep ink black — primary
INK_LIGHT  = RGBColor(0x6B,0x6B,0x6B)  # 淡墨 light ink
SCROLL     = RGBColor(0xF5,0xEC,0xD8)  # 旧宣纸 aged rice paper background
JADE       = RGBColor(0x4A,0x80,0x6E)  # 玉色 jade — nature accent
VERMILLION = RGBColor(0xC0,0x39,0x2B)  # 印章红 seal red — accent / headers / badge

# Variety palette for word cards / games / booklet
MAGENTA = RGBColor(0xD8,0x1B,0x60)
YELLOW  = RGBColor(0xFF,0xC1,0x07)
SKY     = RGBColor(0x42,0xA5,0xF5)
CORAL   = RGBColor(0xFF,0x70,0x43)
PURPLE  = RGBColor(0x9C,0x27,0xB0)
GREEN_L = RGBColor(0x66,0xBB,0x6A)

WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xFF,0xF3,0xE0)
IMGBG   = RGBColor(0xEC,0xE4,0xCF)  # rice-paper-tinted image placeholder
GREEN_OK= RGBColor(0x38,0x8E,0x3C)

CREAM = SCROLL  # alias so D1 helpers using CREAM still produce the scroll background

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=VERMILLION,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n")
    nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph();p.text=line
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    palette=[VERMILLION,JADE,INK,INK_LIGHT,YELLOW,SKY]
    dot_colors=[c for c in palette if c!=color][:4]
    positions=[(0.8,4.7),(1.6,4.5),(7.8,4.5),(8.6,4.7)]
    for (x,y),cl in zip(positions,dot_colors):
        d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(0.4),Inches(0.4))
        d.fill.solid();d.fill.fore_color.rgb=cl;d.line.fill.background()
    return s

def dot(s,x,y,r,color):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(r),Inches(r))
    d.fill.solid();d.fill.fore_color.rgb=color;d.line.fill.background()

def example_slide(cat_emoji, cat_cn, work_cn, work_en, artist, fact, img_lb, color):
    s=ns();bg(s,SCROLL)
    hb(s,f"{cat_emoji} {cat_cn}  ·  《{work_cn}》",color)
    ib(s,0.5,0.95,9,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,work_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,artist,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def example_slide_generic(cat_emoji, cat_cn, name_cn, name_en, sub_label, fact, img_lb, color):
    s=ns();bg(s,SCROLL)
    hb(s,f"{cat_emoji} {cat_cn}  ·  {name_cn}",color)
    ib(s,0.5,0.95,9,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,name_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,sub_label,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def type_overview_slide(emoji, name_cn, name_en, color, hint, subtypes):
    s=ns();bg(s,SCROLL);hb(s,f"{emoji} {name_cn}  {name_en}",color)
    tb(s,0.4,0.85,9,0.35,hint,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    cols = 3 if len(subtypes)<=6 else 4
    rows = (len(subtypes)+cols-1)//cols
    card_w = (9.4 - 0.15*(cols-1))/cols
    card_h = 1.7 if rows<=2 else 1.45
    y_start = 1.3 if rows<=2 else 1.25
    y_step = card_h + 0.25 if rows<=2 else card_h + 0.2
    for i,(sem,scn,sen) in enumerate(subtypes):
        col=i%cols;row=i//cols
        x=0.3+col*(card_w+0.15);y=y_start+row*y_step
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(card_w),Inches(card_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
        tb(s,x+0.1,y+0.1,card_w-0.2,0.55,sem,sz=30,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.75,card_w-0.2,0.4,scn,sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+1.15,card_w-0.2,0.3,sen,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    return s

def word_card_read(w,py,en,sent,img,color=VERMILLION):
    s=ns();bg(s,SCROLL);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=color,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=color)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,img,color=VERMILLION):
    s=ns();bg(s,SCROLL);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=color)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=color)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    return s

n=0

# ============================================================
# 1 COVER — Seal stamp badge with ink-splash dots
# ============================================================
s=ns();n+=1;bg(s,SCROLL)
tb(s,1,0.3,8,0.7,"Chinese Ink Painting",sz=32,b=True,c=INK,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.45,"中国水墨画",sz=22,c=INK_LIGHT,a=PP_ALIGN.CENTER)
# Big square seal-stamp badge (vermillion red, like 印章)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.5),Inches(1.55),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=VERMILLION;sh.line.color.rgb=INK;sh.line.width=Pt(4)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.7),Inches(1.75),Inches(2.6),Inches(2.6))
sh2.fill.solid();sh2.fill.fore_color.rgb=VERMILLION;sh2.line.color.rgb=WHITE;sh2.line.width=Pt(2)
tf=tb(s,3.7,1.85,2.6,0.4,"DAY 3",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"🖌️",sz=52,a=PP_ALIGN.CENTER)
ap(tf,"中国水墨画",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"CHINESE INK",sz=11,c=WARM,a=PP_ALIGN.CENTER)
# ink-splash dots: ink black + light ink + jade
for x,y,r,c in [(1.2,1.9,0.55,INK),(1.5,3.8,0.45,INK_LIGHT),(7.9,1.7,0.5,JADE),(8.2,3.7,0.4,INK),(1.0,4.7,0.32,INK_LIGHT),(8.6,4.6,0.32,JADE),(2.2,2.6,0.22,INK),(7.4,3.0,0.2,INK_LIGHT)]:
    dot(s,x,y,r,c)
tb(s,1,5.05,8,0.4,"🖌️ 一支毛笔，一砚墨——画出中国！One brush, one ink, one painting!",sz=14,b=True,c=VERMILLION,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","认识中国水墨画 + 比较中西画法",VERMILLION),
    ("Session 2  下午","2:00-2:45","复习 + 我会认 5 词 · 我会写 3 词",JADE),
    ("Session 3  下午","3:00-4:30","动手画水墨！基本技巧 + 主题练习 + 自由创作",INK)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=WARM)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"🖌️ 内容目标  Content:",sz=20,b=True,c=VERMILLION)
tf=tb(s,0.7,1.4,9,1.5,"1. 了解中国水墨画的特点（毛笔、墨、水、颜色、浓淡、留白、自然主题）",sz=14,c=DARK)
ap(tf,"2. 认识水墨画常见的主题（竹子、莲花、鱼、熊猫、山水、梅花）",sz=14,c=DARK)
ap(tf,"3. 比较中国水墨画和西洋画的不同",sz=14,c=DARK)
ap(tf,"4. 自己尝试画水墨",sz=14,c=DARK)
tb(s,0.5,3.1,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=JADE)
tb(s,0.7,3.6,9,0.4,"👀 我会认：中国  竹子  花  鱼  熊猫",sz=15,b=True,c=DARK)
tb(s,0.7,4.05,9,0.4,"✍️ 我会写：中国  竹子  花",sz=15,b=True,c=DARK)
tb(s,0.5,4.65,9,0.5,"🖌️ 实践目标：完成 D3 booklet + 1 幅水墨作品",sz=15,c=INK)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","认识中国水墨画 + 比较中西画法\n🖌️ 毛笔 + 墨 + 水 = 水墨画",VERMILLION,"🖌️")
n+=1

# ============================================================
# 5 GALLERY INTRO — preview the 4 examples we'll see one by one
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🖼️ 看一看  Look at These Paintings",VERMILLION)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.3),Inches(8.0),Inches(1.6))
sh.fill.solid();sh.fill.fore_color.rgb=VERMILLION;sh.line.fill.background()
tb(s,1.2,1.45,7.6,0.6,"接下来，我们一起看 4 幅画",sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.2,2.05,7.6,0.5,"Next, we look at 4 paintings together",sz=15,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1.2,2.5,7.6,0.4,"它们都是「中国水墨画」 — 你能发现什么共同点？",sz=14,c=YELLOW,a=PP_ALIGN.CENTER)
# 4 small thumbnails preview
items=[("🎋","《竹子》"),("🦐","《虾》"),("🐼","《熊猫》"),("🏔️","《山水》")]
for i,(em,t) in enumerate(items):
    x=0.6+i*2.3
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.4),Inches(2.0),Inches(1.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=INK;sh.line.width=Pt(2)
    tb(s,x+0.1,3.5,1.8,0.65,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.2,1.8,0.4,t,sz=14,b=True,c=INK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.6,1.8,0.3,str(i+1),sz=10,c=VERMILLION,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.05),Inches(9.4),Inches(0.45))
sh.fill.solid();sh.fill.fore_color.rgb=VERMILLION;sh.line.fill.background()
tb(s,0.4,5.1,9.2,0.4,"👀 仔细看 — 不要急着回答，让眼睛先告诉你！",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 5a-5d  ONE PAINTING PER SLIDE (big image, no features revealed)
# ============================================================
gallery=[
    ("📷 《竹子》郑板桥  Bamboo by Zheng Banqiao","《竹子》郑板桥","🎋",1),
    ("📷 《虾》齐白石  Shrimp by Qi Baishi","《虾》齐白石","🦐",2),
    ("📷 《熊猫》吴作人  Panda by Wu Zuoren (1908-1997)","《熊猫》吴作人","🐼",3),
    ("📷 《溪山行旅图》范宽  Landscape by Fan Kuan","《溪山行旅图》范宽","🏔️",4),
]
GALLERY_NOTES_D3={
    1:"老师备课:\n• 郑板桥 (1693-1765), 清代画家、书法家\n• 「扬州八怪」之一, 一生只画三样: 竹、兰、石\n• 他的竹子瘦、硬、有节 — 像他自己的脾气: 正直\n• 著名诗句: 「咬定青山不放松, 立根原在破岩中」\n• 追问: 「竹子为什么一节一节？」「像不像在长大？」",
    2:"老师备课:\n• 齐白石 (1864-1957), 中国近现代水墨大师\n• 出身木匠, 60 岁后才出名 — 大器晚成\n• 最有名的就是画虾, 看似简单, 但极难画\n• 他研究真虾几十年, 才画出「虾的活」\n• 用浓墨画虾头, 淡墨画身体, 飞白画须 — 一笔一画都有讲究\n• 追问: 「数一数有几只虾？」「它们在做什么？」「为什么没画水？」(留白即水)",
    3:"老师备课:\n• 吴作人 (1908-1997), 中国现代美术大师, 中央美术学院前院长\n• 早年学西洋油画, 后回归中国水墨\n• 他把西画的体积感 + 中国水墨的笔墨结合\n• 最有名的就是画熊猫和金鱼, 1973-1980 年代创作了一系列熊猫作品\n• 用浓墨画耳朵 / 眼圈 / 四肢, 留白画身体 — 黑白对比天然\n• 趣闻: 熊猫的黑白本来就像水墨画! 简直是为水墨画而生的动物\n• 图片参考: 搜索「吴作人 熊猫」可以找到他的代表作\n• 收藏地: 中央美院美术馆, 中国美术馆\n• 追问: 「熊猫为什么是黑白的？」「画家用了多少颜色？」「水墨画里的熊猫和照片里的熊猫有什么不一样？」",
    4:"老师备课:\n• 范宽 (约 950-1032), 北宋画家\n• 《溪山行旅图》是中国十大名画之一, 现存台北故宫博物院\n• 画面气势磅礴: 山很大很高, 旅人和驴子很小很小\n• 这表现了「人在自然中很渺小」的哲学\n• 提示学生: 找一找小小的旅人! (在画中下方, 很难找)\n• 追问: 「你能找到人吗？」「为什么山画得那么大？」「人画得那么小？」",
}
for img_lb,title,em,idx in gallery:
    s=ns();n+=1;bg(s,SCROLL);hb(s,f"🖼️ 看 — 第 {idx} 幅 / {len(gallery)}  Look · {idx}/{len(gallery)}",VERMILLION)
    ib(s,0.5,1.0,9.0,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.6),Inches(9.0),Inches(0.45))
    sh.fill.solid();sh.fill.fore_color.rgb=INK;sh.line.fill.background()
    tb(s,0.6,4.65,1.0,0.35,em,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.6,4.68,7.8,0.35,title,sz=15,b=True,c=WHITE)
    tb(s,0.5,5.15,9.0,0.35,"👀 颜色？主题？空白？— 你看到了什么？",sz=12,c=VERMILLION,a=PP_ALIGN.CENTER)
    pn(s,n)
    if idx in GALLERY_NOTES_D3:
        notes(s,GALLERY_NOTES_D3[idx])

# ============================================================
# 6 INQUIRY — what do you see? (Student observation, no answers)
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🤔 你看到什么？  What Did You See?",VERMILLION)
tb(s,0.4,0.9,9.2,0.35,"看完 4 张画，一起想一想 / After looking, let's think together",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# LEFT: 5 observation questions
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.35),Inches(5.3),Inches(3.7))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=VERMILLION;sh.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.35),Inches(5.3),Inches(0.5))
head.fill.solid();head.fill.fore_color.rgb=VERMILLION;head.line.fill.background()
tb(s,0.45,1.43,5.0,0.4,"❓ 一起想一想  Let's Think",sz=14,b=True,c=WHITE)
qs=[
    "颜色多 还是 少？",
    "主要是什么颜色？",
    "画的是真人，还是自然 (动物 / 植物 / 山)？",
    "画里有没有空白？为什么？",
    "看起来 — 安静 还是 热闹？",
]
tf=tb(s,0.5,2.0,5.0,0.5,f"❓ {qs[0]}",sz=13,c=DARK)
for q in qs[1:]:
    ap(tf,"",sz=6)
    ap(tf,f"❓ {q}",sz=13,c=DARK)
# RIGHT: sentence frames
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.85),Inches(1.35),Inches(3.85),Inches(3.7))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=JADE;sh.line.width=Pt(2.5)
head2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.85),Inches(1.35),Inches(3.85),Inches(0.5))
head2.fill.solid();head2.fill.fore_color.rgb=JADE;head2.line.fill.background()
tb(s,6.0,1.43,3.55,0.4,"💬 我会说  I Can Say",sz=14,b=True,c=WHITE)
frames=[
    "我看到 ________。",
    "颜色 ________。",
    "画的是 ________。",
    "画里 ________。",
    "我觉得 ________。",
]
tf2=tb(s,6.05,2.0,3.65,0.5,f"·  {frames[0]}",sz=13,c=DARK)
for fr in frames[1:]:
    ap(tf2,"",sz=6)
    ap(tf2,f"·  {fr}",sz=13,c=DARK)
# Bottom transition
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.15),Inches(9.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=INK;sh.line.fill.background()
tb(s,0.4,5.18,9.2,0.35,"👉 大家说完，我们一起总结 — 这就是「水墨画的特点」!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 7 REVEAL — big idea (now confirms what students discovered)
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"💡 你说对啦！水墨画 = ?",VERMILLION)
tb(s,0.4,0.9,9.2,0.35,"You discovered it! Here's the formula.",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.4),Inches(8.0),Inches(1.5))
sh.fill.solid();sh.fill.fore_color.rgb=INK;sh.line.fill.background()
tb(s,1.2,1.5,7.6,0.7,"毛笔 + 墨 + 水 = 水墨画",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.2,2.2,7.6,0.55,"Brush + Ink + Water = Chinese Ink Painting",sz=18,c=WARM,a=PP_ALIGN.CENTER)
bubbles=[
    ("🖌️","工具","Tools",INK),
    ("⚫","颜色","Colors",INK_LIGHT),
    ("🍃","感觉","Feeling",JADE),
]
for i,(em,cn,en,cl) in enumerate(bubbles):
    x=0.5+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.2),Inches(3.0),Inches(1.95))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,3.3,2.8,0.65,em,sz=40,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.0,2.8,0.45,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.55,2.8,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8 SIX FEATURES — 2x3 grid (summary of what students discovered)
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"📌 总结  6 Features of Ink Painting")
tb(s,0.4,0.85,9.2,0.3,"我们一起发现的 6 个特点 / The 6 features we found together",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
features=[
    ("🖌️","工具","Brush, ink, water, rice paper\n毛笔、墨、水、宣纸",INK),
    ("🎨","颜色","Mostly black/white/gray\n黑白灰为主",INK_LIGHT),
    ("💧","浓淡","Dark or light\n深一点 / 浅一点",JADE),
    ("🌿","主题","Nature: bamboo, fish, panda…\n自然 — 竹子、花、鱼、熊猫",GREEN_L),
    ("🌬️","留白","Don't fill it up — leave space\n不用画满",VERMILLION),
    ("🧘","感觉","Quiet, simple, poetic\n安静、简单、有意境",PURPLE),
]
for i,(em,cn,en,cl) in enumerate(features):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.2+row*2.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.1,2.8,0.65,em,sz=34,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.8,2.8,0.4,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    ls=en.split('\n')
    tb(s,x+0.1,y+1.25,2.8,0.3,ls[0],sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    if len(ls)>1:
        tb(s,x+0.1,y+1.5,2.8,0.3,ls[1],sz=10,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 7 浓墨 vs 淡墨 demo
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"💧 浓墨 vs 淡墨  Dark Ink vs Light Ink",JADE)
tb(s,0.4,0.9,9.2,0.4,"水多 = 浅墨 (淡)  ·  水少 = 浓墨 (深)",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
# Left: dark ink card
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.4),Inches(4.5),Inches(3.2))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=INK;sh.line.width=Pt(3)
tb(s,0.5,1.5,4.3,0.5,"🖤 浓墨  Dark Ink",sz=22,b=True,c=INK,a=PP_ALIGN.CENTER)
# A heavy dark stroke (rounded rectangle)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.2),Inches(2.2),Inches(2.9),Inches(0.5))
sh2.fill.solid();sh2.fill.fore_color.rgb=INK;sh2.line.fill.background()
tb(s,0.5,2.85,4.3,0.4,"水少 + 多墨",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,3.25,4.3,0.4,"Less water + more ink",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.5,3.7,4.3,0.4,"用来画：树枝、眼睛、轮廓",sz=12,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,4.1,4.3,0.4,"For: branches, eyes, outlines",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Right: light ink card
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.4),Inches(4.5),Inches(3.2))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=INK_LIGHT;sh.line.width=Pt(3)
tb(s,5.2,1.5,4.3,0.5,"💧 淡墨  Light Ink",sz=22,b=True,c=INK_LIGHT,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(2.2),Inches(2.9),Inches(0.5))
sh2.fill.solid();sh2.fill.fore_color.rgb=INK_LIGHT;sh2.line.fill.background()
tb(s,5.2,2.85,4.3,0.4,"水多 + 少墨",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,5.2,3.25,4.3,0.4,"More water + less ink",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,5.2,3.7,4.3,0.4,"用来画：山、雾、远的东西",sz=12,c=DARK,a=PP_ALIGN.CENTER)
tb(s,5.2,4.1,4.3,0.4,"For: mountains, mist, distance",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Inquiry strip
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.75),Inches(9.2),Inches(0.5))
sh3.fill.solid();sh3.fill.fore_color.rgb=VERMILLION;sh3.line.fill.background()
tb(s,0.5,4.8,9.0,0.4,"🤔 你想试试哪一种先？  Which one will you try first?",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8 留白 inquiry
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🌬️ 留白 — 不用画满  Leave Space — Don't Fill It Up",JADE)
tb(s,0.4,0.9,9.2,0.4,"中国水墨画喜欢留白，让画呼吸。",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
# Left thumbnail: 画满
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(1.4),Inches(4.2),Inches(2.8))
sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.color.rgb=INK;sh.line.width=Pt(2)
tb(s,0.7,2.4,4.0,0.7,"📷 画满的画",sz=20,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.7,3.1,4.0,0.4,"(A busy painting)",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.6,4.25,4.2,0.4,"画满 — 看起来很热闹",sz=14,b=True,c=INK,a=PP_ALIGN.CENTER)
tb(s,0.6,4.65,4.2,0.3,"Filled up — busy, lively",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Right thumbnail: 留白
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.4),Inches(4.2),Inches(2.8))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=JADE;sh.line.width=Pt(2)
# small ink mark in corner to suggest sparseness
sh_m=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(7.0),Inches(2.5),Inches(0.4),Inches(0.4))
sh_m.fill.solid();sh_m.fill.fore_color.rgb=INK;sh_m.line.fill.background()
tb(s,5.3,3.1,4.0,0.4,"📷 留白的画 (mostly empty)",sz=12,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,5.2,4.25,4.2,0.4,"留白 — 看起来很安静",sz=14,b=True,c=JADE,a=PP_ALIGN.CENTER)
tb(s,5.2,4.65,4.2,0.3,"With empty space — quiet, peaceful",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Inquiry strip
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.15),Inches(9.2),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=VERMILLION;sh3.line.fill.background()
tb(s,0.5,5.18,9.0,0.35,"🤔 哪一幅看起来更安静？为什么？",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 9 COMMON SUBJECTS — type_overview
# ============================================================
s=type_overview_slide("🌿","水墨画的主题","Common Subjects",JADE,
    "水墨画喜欢画什么？大多是大自然！",
    [("🎋","竹子","Bamboo"),("🪷","莲花","Lotus"),
     ("🐟","鱼","Fish"),("🐼","熊猫","Panda"),
     ("⛰️","山水","Landscape"),("🌸","梅花","Plum Blossom")]);n+=1;pn(s,n)

# ============================================================
# 10 VIDEO — see ink painting in action (replaces duplicate bamboo)
# ============================================================
s=ns();n+=1;bg(s,INK)
tb(s,1,0.7,8,0.8,"🎬 看视频  Watch a Video",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.55,8,0.45,"水墨画家是怎么画的？  How does an ink painter paint?",sz=18,c=WARM,a=PP_ALIGN.CENTER)
ib(s,1.2,2.2,7.6,2.5,"📷 视频截图位置 / Video screenshot")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.85),Inches(9),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=VERMILLION;sh.line.fill.background()
tb(s,0.6,4.9,8.8,0.4,"🔗 推荐: YouTube 搜「Chinese ink painting demo」/ B站 搜「水墨画 教学」",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)
notes(s,"老师备课:\n• 选 1-2 段 2-3 分钟的水墨画过程视频, 让学生看到「真正在画」\n• 重点观察: 笔的姿势 / 蘸墨蘸水的多少 / 一笔下去的速度\n• 推荐搜索:\n   - YouTube: \"Chinese ink painting demo\" / \"sumi-e brush demo\"\n   - 哔哩哔哩 (B站): \"水墨画 教程\" / \"齐白石 画虾 过程\"\n   - 抖音/小红书: \"国画 教学\"\n• 看完追问: 「画家先画哪里？」「水多还是水少？」「他犹豫了吗？」\n• 趣闻: 水墨画一笔下去不能改, 所以画家要先在心里「画」一遍")

# ============================================================
# 11 FISH example — 八大山人 (different artist, covers 鱼 language goal)
# ============================================================
s=example_slide("🐟","鱼 Fish","鳜鱼图","Mandarin Fish",
    "八大山人 Bada Shanren · 朱耷 · 1626-1705 · 明末清初",
    "一条鱼 + 一片空白 = 一个世界。\nOne fish + lots of empty space = a whole world!",
    "📷 八大山人《鱼》(搜「八大山人 鱼」)",JADE);n+=1;pn(s,n)
notes(s,"老师备课:\n• 八大山人 (本名朱耷, 1626-1705) 明朝皇室后裔, 明亡后出家当和尚\n• 他的画很特别: 鱼的眼睛是「白眼向天」— 表达对当时世界的不满\n• 画一条孤独的鱼, 周围全是空白 — 这就是「极致留白」\n• 他常常只画一个东西: 一条鱼、一只鸟、一朵荷花\n• 简到不能再简, 但每一笔都很有力量\n• 图片参考: 网搜「八大山人 鱼」/「朱耷 鱼」\n• 收藏: 北京故宫博物院 / 上海博物馆\n• 追问: 「为什么只画一条鱼？」「鱼的眼睛在哪里？为什么这么画？」「这条鱼开心还是难过？」")

# ============================================================
# 12 LOTUS example — 张大千 (covers 花 language goal, new subject!)
# ============================================================
s=example_slide("🪷","花 Flower (荷花)","泼墨荷花","Splashed Ink Lotus",
    "张大千 Zhang Daqian · 1899-1983 · 中国现代水墨大师",
    "「泼墨」— 把墨泼上去！花就出来了！\nSplash the ink — and a flower appears!",
    "📷 张大千《荷花》(搜「张大千 荷花 泼墨」)",JADE);n+=1;pn(s,n)
notes(s,"老师备课:\n• 张大千 (1899-1983) 是 20 世纪最有名的中国画家之一\n• 他什么都画 — 山水、花鸟、人物, 还会临摹古画 (传说曾骗过专家)\n• 后来发明了「泼墨泼彩」技法 — 把墨和颜料直接泼到纸上\n• 这种画法很大胆, 一旦泼下去就改不了\n• 荷花是中国的传统题材 (出淤泥而不染), 但张大千的荷花是「现代的」\n• 图片参考: 网搜「张大千 荷花」/「张大千 泼墨」\n• 收藏: 台北故宫博物院 / 苏富比拍卖纪录\n• 趣闻: 他的画是当今拍卖纪录最高的中国画家之一, 一幅荷花拍过几亿\n• 追问: 「「泼」墨是什么意思？」「荷花和别的花有什么不一样？」「为什么古人爱画荷花？」")

# ============================================================
# 13 PANDA example — 韩美林 (different from gallery's 吴作人)
# ============================================================
s=example_slide_generic("🐼","熊猫 Panda","卡通熊猫","Folk-Style Panda",
    "韩美林 Han Meilin · 1936-至今 · 现代艺术家",
    "他的熊猫圆滚滚 — 像不像奥运福娃？\nHis pandas are round and cute — like Olympic mascots!",
    "📷 韩美林《熊猫》(搜「韩美林 熊猫」)",JADE);n+=1;pn(s,n)
notes(s,"老师备课:\n• 韩美林 (1936-至今) 中国现代艺术家、设计师\n• 他设计了 2008 北京奥运会吉祥物「福娃」, 也是熊猫主题设计大师\n• 他的熊猫不是「写实」的, 是「卡通化」的水墨 — 圆圆的, 萌萌的\n• 用最简单的几笔, 就抓住熊猫的特征\n• 这种风格融合了中国民间艺术 + 现代设计\n• 图片参考: 网搜「韩美林 熊猫」\n• 趣闻: 他设计的福娃熊猫 (晶晶) 全世界 30 亿人都认识\n• 收藏: 韩美林艺术馆 (北京 / 杭州 / 银川)\n• 对比: 吴作人画熊猫像写真, 韩美林画熊猫像卡通\n• 追问: 「这只熊猫和上一幅 (吴作人) 不一样在哪里？」「你更喜欢哪一种？为什么？」")

# ============================================================
# 14 比较 中国水墨画 vs 西洋画
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🆚 中国水墨画 vs 西洋画  Chinese Ink vs Western Painting",VERMILLION)
ts=s.shapes.add_table(6,3,Inches(0.3),Inches(0.95),Inches(9.4),Inches(4.3));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(3.8);t.columns[2].width=Inches(3.8)
rows=[
    ["",     "🇨🇳 中国水墨画",                "🇪🇺 西洋画"],
    ["颜色", "颜色少 (黑白灰)\nFew colors (B/W/gray)",      "颜色多\nMany colors"],
    ["像不像真","重感觉，不一定像\nFeeling > realism",        "像照片 一样真\nLooks like a photo"],
    ["留白", "讲究留白\nLeaves empty space",                  "画得满\nFills the canvas"],
    ["感觉", "安静、简单\nQuiet, simple",                     "颜色丰富、细节多\nRich, detailed"],
    ["例子", "齐白石 · 虾\nQi Baishi — Shrimp",               "梵高 · 向日葵\nVan Gogh — Sunflowers"],
]
for r,rd in enumerate(rows):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True
        p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(12 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        if r==0:
            rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=VERMILLION
        elif c==0:
            rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        elif c==1:
            rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xFA,0xF4,0xE5)
        else:
            rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xEC,0xF2,0xEF)
pn(s,n)

# ============================================================
# 15 PAIR-UP inquiry — sort 中 vs 西
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🤔 找一找  Sort the Paintings",VERMILLION)
tb(s,0.4,0.9,9.2,0.4,"看 4 幅画 — 哪些是中国水墨画？哪些是西洋画？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
pics=[
    ("📷 齐白石 虾","Qi Baishi · Shrimp","🇨🇳 水墨",JADE),
    ("📷 梵高 向日葵","Van Gogh · Sunflowers","🇪🇺 西洋",MAGENTA),
    ("📷 范宽 山水","Fan Kuan · Landscape","🇨🇳 水墨",JADE),
    ("📷 蒙娜丽莎","Mona Lisa","🇪🇺 西洋",MAGENTA),
]
for i,(img,cap,tag,cl) in enumerate(pics):
    col=i%4
    x=0.3+col*2.4;y=1.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.2),Inches(3.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    ib(s,x+0.1,y+0.15,2.0,1.7,img)
    tb(s,x+0.1,y+1.95,2.0,0.4,cap,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(y+2.45),Inches(1.8),Inches(0.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+0.2,y+2.5,1.8,0.4,tag,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.4,4.85,9.2,0.4,"🗣️ 说说看：你怎么知道？(颜色？留白？感觉？)",sz=13,b=True,c=VERMILLION,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 16 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","复习 + 语言目标\n👀 我会认 5 个词  ·  ✍️ 我会写 3 个词",JADE,"📖")
n+=1

# ============================================================
# 17 REVIEW match — 主题 ↔ feature
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🔄 快速复习  Quick Review — 连一连",JADE)
tb(s,0.4,0.85,9,0.35,"把主题和它的特点连起来 (口头)",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
subjects=[("🎋 竹子",JADE),("🌸 花",MAGENTA),("🐟 鱼",SKY),("🐼 熊猫",INK)]
hints=["黑白颜色刚刚好","节节高，画很多线条","水里游，齐白石画得最有名","花瓣淡淡的，留一点白"]
for i,(em,cl) in enumerate(subjects):
    y=1.25+i*0.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(y),Inches(3.4),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.12,3.2,0.5,em,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
for i,h in enumerate(hints):
    y=1.25+i*0.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.6),Inches(y),Inches(3.8),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=JADE;sh.line.width=Pt(2)
    tb(s,5.7,y+0.15,3.6,0.5,h,sz=14,c=DARK,a=PP_ALIGN.CENTER)
tb(s,4.2,3.0,1.6,0.6,"?",sz=44,b=True,c=VERMILLION,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 18-22 我会认 word cards
# ============================================================
read_words=[
    ("中国","zhōng guó","China","中国水墨画很有名。","📷 中国地图 / 国旗",VERMILLION),
    ("竹子","zhú zi","bamboo","竹子是水墨画常画的。","📷 竹子",JADE),
    ("花","huā","flower","我画了一朵花。","📷 一朵花",MAGENTA),
    ("鱼","yú","fish","齐白石画的鱼很活。","📷 水墨鱼",SKY),
    ("熊猫","xióng māo","panda","熊猫是中国的国宝。","📷 熊猫",INK),
]
for w,py,en,sent,img,cl in read_words:
    s=word_card_read(w,py,en,sent,img,cl);n+=1;pn(s,n)

# ============================================================
# 23 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🎮 练一练  Word Games (选一个玩！)",JADE)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","老师说词\n学生拍正确的字卡",RGBColor(0xFF,0xF3,0xE0)),
    ("2️⃣","配主题图\nMatch Picture","看水墨画\n找对应的字卡",RGBColor(0xE3,0xF2,0xFD)),
    ("3️⃣","字卡接龙\nWord Chain","一个接一个\n用字卡造句",RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣","我画你猜\nDraw & Guess","小朋友画\n大家猜词",RGBColor(0xFC,0xE4,0xEC)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 24-26 我会写 word cards
# ============================================================
write_words=[
    ("中国","zhōng guó","China","📷 中国国旗 / 地图",VERMILLION),
    ("竹子","zhú zi","bamboo","📷 竹子",JADE),
    ("花","huā","flower","📷 一朵花",MAGENTA),
]
for w,py,en,img,cl in write_words:
    s=word_card_write(w,py,en,img,cl);n+=1;pn(s,n)

# ============================================================
# 27 SESSION 3 DIVIDER
# ============================================================
div("Session 3  下午","动手画水墨！  Make Ink Art!\n🖌️ 基本技巧 + 主题练习 + 自由创作",INK,"🖌️")
n+=1

# ============================================================
# 28 BOOKLET
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"📓 完成 D3 练习册  Day 3 Booklet",INK)
ib(s,0.4,0.9,9.2,4.3,"📷 D3 练习册截图  (teacher: insert booklet pages)")
pn(s,n)

# ============================================================
# 29 PROJECT OVERVIEW — 3 steps
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🖌️ 今天怎么画？  Today's 3 Steps",INK)
tb(s,0.4,0.9,9,0.4,"一步一步来 — 每一步都有练习！Step by step.",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("STEP 1","🖌️ 练习基本技巧","Basic Brush Skills","浓墨 / 淡墨 / 调水\n线条 / 点",VERMILLION),
    ("STEP 2","🌿 练习常见主题","Common Subjects","竹子 / 莲花\n鱼/虾 / 熊猫",JADE),
    ("STEP 3","🎨 自由创作","Free Creation","选一个主题\n完成你的作品",INK),
]
for i,(num,cn,en,d,cl) in enumerate(steps):
    x=0.4+i*3.1
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(2.9),Inches(3.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3.5)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(1.55),Inches(1.4),Inches(0.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+0.25,1.6,1.3,0.4,num,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.15,2.7,0.6,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.75,2.7,0.4,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,3.4,2.7,0.4,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=13,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 30 STEP 1 — 基本技巧 details
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🖌️ Step 1: 基本技巧  Basic Brush Skills",VERMILLION)
# Materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=VERMILLION;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"🖌️ 毛笔  Brush",sz=13,c=DARK)
ap(tf,"🖤 墨 (墨汁 / 墨条 + 砚)  Ink",sz=13,c=DARK)
ap(tf,"💧 水盘  Water dish",sz=13,c=DARK)
ap(tf,"📄 宣纸 / 吸水纸  Rice paper / absorbent paper",sz=13,c=DARK)
# Tip strip
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=YELLOW;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"💡 小贴士  Tip",sz=14,b=True,c=DARK)
tf_k=tb(s,0.4,3.85,4.4,1.4,"· 笔尖蘸墨，笔肚蘸水",sz=13,c=DARK)
ap(tf_k,"· 多试一试 — 没有对错！",sz=13,c=DARK)
ap(tf_k,"· 慢慢画 — 毛笔会跟你说话",sz=13,c=DARK)
# Right: 4 mini-exercises
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=INK;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"✏️ 4 个小练习  4 Mini-Exercises",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,3.6,"1️⃣ 浓墨：直接蘸墨画一笔",sz=13,b=True,c=INK)
ap(tf2,"   Dip in ink, draw one bold line",sz=10,c=GRAY)
ap(tf2,"2️⃣ 淡墨：先蘸水，再蘸一点点墨",sz=13,b=True,c=INK_LIGHT)
ap(tf2,"   Wet first, then a touch of ink",sz=10,c=GRAY)
ap(tf2,"3️⃣ 浓变淡：一笔从浓到淡 (\"飞白\")",sz=13,b=True,c=VERMILLION)
ap(tf2,"   One stroke: dark → light (\"flying white\")",sz=10,c=GRAY)
ap(tf2,"4️⃣ 点 + 线：试试不同压力",sz=13,b=True,c=JADE)
ap(tf2,"   Dots & lines: try different pressure",sz=10,c=GRAY)
pn(s,n)

# ============================================================
# 31 STEP 2 — 主题练习 (竹子 + 鱼)
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🌿 Step 2: 主题练习 (竹子 + 鱼)  Bamboo + Fish",JADE)
# Left: 竹子
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.6),Inches(4.3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=JADE;sh.line.width=Pt(3)
tb(s,0.4,1.05,4.4,0.5,"🎋 竹子  Bamboo",sz=22,b=True,c=JADE,a=PP_ALIGN.CENTER)
ib(s,0.5,1.6,4.2,1.4,"📷 竹子分步图")
tf=tb(s,0.5,3.1,4.2,2.1,"1️⃣ 画竹竿 — 一笔直上",sz=13,b=True,c=DARK)
ap(tf,"   Bamboo stalk — one straight stroke",sz=10,c=GRAY)
ap(tf,"2️⃣ 画竹节 — 短短的横线",sz=13,b=True,c=DARK)
ap(tf,"   Joints — short horizontal marks",sz=10,c=GRAY)
ap(tf,"3️⃣ 画竹叶 — 一笔一片",sz=13,b=True,c=DARK)
ap(tf,"   Leaves — one stroke each",sz=10,c=GRAY)
ap(tf,"4️⃣ 加淡墨 — 远的叶子",sz=13,b=True,c=DARK)
ap(tf,"   Light ink — distant leaves",sz=10,c=GRAY)
# Right: 鱼
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(0.95),Inches(4.6),Inches(4.3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SKY;sh.line.width=Pt(3)
tb(s,5.2,1.05,4.4,0.5,"🐟 鱼  Fish",sz=22,b=True,c=SKY,a=PP_ALIGN.CENTER)
ib(s,5.3,1.6,4.2,1.4,"📷 鱼分步图")
tf=tb(s,5.3,3.1,4.2,2.1,"1️⃣ 画身体 — 一个椭圆",sz=13,b=True,c=DARK)
ap(tf,"   Body — an oval",sz=10,c=GRAY)
ap(tf,"2️⃣ 画头 — 加一点小尖",sz=13,b=True,c=DARK)
ap(tf,"   Head — add a little point",sz=10,c=GRAY)
ap(tf,"3️⃣ 画尾 — 像扇子一样",sz=13,b=True,c=DARK)
ap(tf,"   Tail — like a fan",sz=10,c=GRAY)
ap(tf,"4️⃣ 加眼睛 — 一个小黑点",sz=13,b=True,c=DARK)
ap(tf,"   Eye — one little black dot",sz=10,c=GRAY)
pn(s,n)

# ============================================================
# 32 STEP 2 — 主题练习 (莲花 + 熊猫)
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🌿 Step 2: 主题练习 (莲花 + 熊猫)  Lotus + Panda",JADE)
# Left: 莲花
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.6),Inches(4.3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MAGENTA;sh.line.width=Pt(3)
tb(s,0.4,1.05,4.4,0.5,"🪷 莲花  Lotus",sz=22,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ib(s,0.5,1.6,4.2,1.4,"📷 莲花分步图")
tf=tb(s,0.5,3.1,4.2,2.1,"1️⃣ 淡墨花瓣 — 一瓣一瓣",sz=13,b=True,c=DARK)
ap(tf,"   Light-ink petals — one by one",sz=10,c=GRAY)
ap(tf,"2️⃣ 中心圆 — 花心",sz=13,b=True,c=DARK)
ap(tf,"   Center circle — the heart",sz=10,c=GRAY)
ap(tf,"3️⃣ 茎 — 一笔直直的",sz=13,b=True,c=DARK)
ap(tf,"   Stem — one straight stroke",sz=10,c=GRAY)
ap(tf,"4️⃣ 一片荷叶 — 大大的圆",sz=13,b=True,c=DARK)
ap(tf,"   One lotus leaf — a big round shape",sz=10,c=GRAY)
# Right: 熊猫
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(0.95),Inches(4.6),Inches(4.3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=INK;sh.line.width=Pt(3)
tb(s,5.2,1.05,4.4,0.5,"🐼 熊猫  Panda",sz=22,b=True,c=INK,a=PP_ALIGN.CENTER)
ib(s,5.3,1.6,4.2,1.4,"📷 熊猫分步图")
tf=tb(s,5.3,3.1,4.2,2.1,"1️⃣ 淡墨身体 — 一团圆",sz=13,b=True,c=DARK)
ap(tf,"   Light-ink body — a round shape",sz=10,c=GRAY)
ap(tf,"2️⃣ 黑墨耳朵 + 眼圈",sz=13,b=True,c=DARK)
ap(tf,"   Dark ink — ears & eye patches",sz=10,c=GRAY)
ap(tf,"3️⃣ 四肢 — 黑墨四块",sz=13,b=True,c=DARK)
ap(tf,"   Legs — four dark patches",sz=10,c=GRAY)
ap(tf,"4️⃣ 背景一支竹  Add one bamboo",sz=13,b=True,c=DARK)
ap(tf,"   Background — one stalk of bamboo",sz=10,c=GRAY)
pn(s,n)

# ============================================================
# 33 STEP 3 — 自由创作
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🎨 Step 3: 自由创作  Free Creation",INK)
# Top: instruction card
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(9.4),Inches(0.85))
sh.fill.solid();sh.fill.fore_color.rgb=INK;sh.line.fill.background()
tb(s,0.5,1.0,9.0,0.4,"选一个主题，完成你自己的水墨作品！",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,1.4,9.0,0.4,"Pick one subject and finish your own ink painting!",sz=12,c=WARM,a=PP_ALIGN.CENTER)
# 4 subject choice cards
choices=[
    ("🎋","竹子","Bamboo",JADE),
    ("🪷","莲花","Lotus",MAGENTA),
    ("🐟","鱼/虾","Fish/Shrimp",SKY),
    ("🐼","熊猫","Panda",INK),
]
for i,(em,cn,en,cl) in enumerate(choices):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(1.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.1,2.0,0.7,em,sz=44,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.95,2.0,0.4,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.4,2.0,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Sentence frames
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.45))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=VERMILLION;sh3.line.width=Pt(2)
tb(s,0.5,3.95,9.0,0.4,"🗣️ 展示句型  Show & Tell",sz=14,b=True,c=VERMILLION)
tf=tb(s,0.5,4.35,9.0,0.4,"· 这是我画的____。  This is my ____.",sz=13,c=DARK)
ap(tf,"· 我用了浓墨/淡墨。  I used dark / light ink.",sz=13,c=DARK)
ap(tf,"· 我留了一点白。  I left some white space.",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 34 GALLERY WALK
# ============================================================
s=ns();n+=1;bg(s,SCROLL);hb(s,"🖼️ 小小水墨展  Ink Gallery Walk",VERMILLION)
tb(s,0.4,0.9,9,0.5,"把作品挂起来，像博物馆一样看一看！",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Hang your art and walk around like a museum!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("1️⃣","🖼️","挂作品\nHang it up"),
    ("2️⃣","👀","看一看\nLook around"),
    ("3️⃣","🗣️","说一说\nTalk about it"),
    ("4️⃣","⭐","投票最喜欢\nVote favorite"),
]
colors=[VERMILLION,JADE,INK,YELLOW]
for i,(num,em,d) in enumerate(steps):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=colors[i];sh.line.width=Pt(3)
    tb(s,x+0.1,2.1,2.0,0.5,num,sz=22,b=True,c=colors[i],a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.6,2.0,0.8,em,sz=40,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,3.6,2.0,0.4,ls[0],sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 35 DAY 3 BADGE — Vermillion seal stamp
# ============================================================
s=ns();n+=1;bg(s,SCROLL)
tb(s,0.5,0.4,9,0.8,"🎖️ Day 3 小水墨家徽章  Ink Artist Seal",sz=26,b=True,c=VERMILLION,a=PP_ALIGN.CENTER)
# Big square seal stamp
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=VERMILLION;sh.line.color.rgb=INK;sh.line.width=Pt(5)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.7),Inches(1.6),Inches(2.6),Inches(2.6))
sh2.fill.solid();sh2.fill.fore_color.rgb=VERMILLION;sh2.line.color.rgb=WHITE;sh2.line.width=Pt(2)
tf=tb(s,3.7,1.7,2.6,0.4,"DAY 3",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"🖌️",sz=42,a=PP_ALIGN.CENTER)
ap(tf,"中国水墨画",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=12,b=True,c=WARM,a=PP_ALIGN.CENTER)
# Ink-splash dots around
for x,y,cl in [(2.5,2.2,INK),(7.0,2.2,JADE),(2.5,3.7,INK_LIGHT),(7.0,3.7,INK),(2.0,3.0,JADE),(7.5,3.0,INK_LIGHT)]:
    dot(s,x,y,0.32,cl)
tb(s,1,4.55,8,0.4,"恭喜你完成 Day 3！You are an ink artist! 🎉",sz=16,b=True,c=VERMILLION,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"认识水墨画 6 大特点 · 5 个新词 · 1 幅水墨作品",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 36 TOMORROW PREVIEW
# ============================================================
s=ns();n+=1;bg(s,VERMILLION)
tb(s,0.5,0.9,9,0.8,"🎨 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.5,"Day 4 — 当代艺术",sz=28,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
ap(tf,"Contemporary Art",sz=16,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🎨 大胆表达！艺术可以很奇怪很大胆",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Be bold — art can be weird & wild!",sz=14,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，小艺术家！",sz=15,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/小小艺术家little_artist_pbl/day3_chinese_ink.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
