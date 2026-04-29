#!/usr/bin/env python3
"""
小小艺术家 Little Artist Unit — Day 2: 跟着大师学画画 Follow the Masters
Picasso · Matisse · Van Gogh
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

MAGENTA = RGBColor(0xD8,0x1B,0x60)
YELLOW  = RGBColor(0xFF,0xC1,0x07)
SKY     = RGBColor(0x42,0xA5,0xF5)
CORAL   = RGBColor(0xFF,0x70,0x43)
PURPLE  = RGBColor(0x9C,0x27,0xB0)
GREEN_L = RGBColor(0x66,0xBB,0x6A)
CREAM   = RGBColor(0xFF,0xF8,0xE7)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xFF,0xF3,0xE0)
IMGBG   = RGBColor(0xF0,0xE8,0xF0)
GREEN_OK= RGBColor(0x38,0x8E,0x3C)

PICASSO = RGBColor(0x42,0x5B,0x8A)
MATISSE = RGBColor(0xE5,0x39,0x35)
VANGOGH = RGBColor(0xFD,0xB9,0x18)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
import os as _os
_IMG_DIR=_os.path.join(_os.path.dirname(_os.path.abspath(__file__)),"images")
def ib(s,l,t,w,h,lb="📷",url=None,key=None):
    """Image placeholder. If `key` matches a file in images/, embed it directly.
    Otherwise draw the gray box + label + optional URL line."""
    if key:
        for ext in ("jpg","jpeg","png","webp"):
            p=_os.path.join(_IMG_DIR,f"{key}.{ext}")
            if _os.path.exists(p):
                s.shapes.add_picture(p,Inches(l),Inches(t),Inches(w),Inches(h))
                return
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    if url:
        tb(s,l+0.1,t+h/2-0.4,w-0.2,0.35,lb,sz=13,c=LGRAY,a=PP_ALIGN.CENTER)
        tb(s,l+0.1,t+h/2+0.0,w-0.2,0.3,f"🔗 {url}",sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
    else:
        tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=MAGENTA,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n")
    nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph();p.text=line
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    palette=[YELLOW,CORAL,SKY,GREEN_L,PURPLE,MAGENTA,PICASSO,MATISSE,VANGOGH]
    dot_colors=[c for c in palette if c!=color][:4]
    positions=[(0.8,4.7),(1.6,4.5),(7.8,4.5),(8.6,4.7)]
    for (x,y),cl in zip(positions,dot_colors):
        d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(0.4),Inches(0.4))
        d.fill.solid();d.fill.fore_color.rgb=cl;d.line.fill.background()
    return s

def dot(s,x,y,r,color):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(r),Inches(r))
    d.fill.solid();d.fill.fore_color.rgb=color;d.line.fill.background()

def observe_slide(emoji, work_cn, work_en, color, img_lb, questions, url=None, key=None):
    """Observe-first slide: students study painting + answer guided questions BEFORE teacher explains."""
    s=ns();bg(s,CREAM)
    hb(s,f"👀 先看一看 · 《{work_cn}》",color)
    tb(s,0.4,0.9,9.2,0.3,f"Look First · {work_en} — what do you see?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    # LEFT: image
    ib(s,0.3,1.3,4.4,3.6,img_lb,url=url,key=key)
    # RIGHT: 4 questions panel
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.3),Inches(4.85),Inches(3.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.3),Inches(4.85),Inches(0.5))
    head.fill.solid();head.fill.fore_color.rgb=color;head.line.fill.background()
    tb(s,5.0,1.38,4.6,0.4,"🤔 你看到什么？  What Do You See?",sz=14,b=True,c=WHITE)
    tf=tb(s,5.05,1.95,4.55,0.5,f"❓ {questions[0]}",sz=13,c=DARK)
    for q in questions[1:]:
        ap(tf,"",sz=6)
        ap(tf,f"❓ {q}",sz=13,c=DARK)
    # bottom hint
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.0),Inches(9.4),Inches(0.4))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.4,5.04,9.2,0.32,"👁️ 不要急 — 让眼睛先告诉你! 然后听老师讲。",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    return s

def example_slide(cat_emoji, cat_cn, work_cn, work_en, artist, fact, img_lb, color, url=None, key=None, trivia=None, connect=None):
    s=ns();bg(s,CREAM)
    hb(s,f"{cat_emoji} {cat_cn}  ·  《{work_cn}》",color)
    if trivia or connect:
        # Rich layout: smaller image + bigger info card with 趣闻 / 你在哪里看过
        ib(s,0.5,0.95,9,2.85,img_lb,url=url,key=key)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.9),Inches(9),Inches(1.5))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
        tb(s,0.65,3.95,4.5,0.32,work_en,sz=14,b=True,c=color)
        tb(s,5.0,3.97,4.3,0.28,artist,sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
        first_fact=fact.split("\n")[0]
        tb(s,0.65,4.30,8.7,0.32,f"💬 {first_fact}",sz=12,c=DARK)
        if trivia:
            tb(s,0.65,4.62,8.7,0.32,f"💡 趣闻: {trivia}",sz=11,b=True,c=color)
        if connect:
            tb(s,0.65,4.97,8.7,0.32,f"👀 {connect}",sz=11,c=DARK)
    else:
        ib(s,0.5,0.95,9,3.5,img_lb,url=url,key=key)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
        tb(s,0.7,4.6,4.5,0.4,work_en,sz=15,b=True,c=color)
        tb(s,0.7,4.95,4.5,0.3,artist,sz=11,c=GRAY)
        tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def example_slide_generic(cat_emoji, cat_cn, name_cn, name_en, sub_label, fact, img_lb, color, url=None, key=None):
    s=ns();bg(s,CREAM)
    hb(s,f"{cat_emoji} {cat_cn}  ·  {name_cn}",color)
    ib(s,0.5,0.95,9,3.5,img_lb,url=url,key=key)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,name_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,sub_label,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def type_overview_slide(emoji, name_cn, name_en, color, hint, subtypes):
    s=ns();bg(s,CREAM);hb(s,f"{emoji} {name_cn}  {name_en}",color)
    tb(s,0.4,0.85,9,0.35,hint,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    cols = 3 if len(subtypes)<=6 else 4
    rows = (len(subtypes)+cols-1)//cols
    card_w = (9.4 - 0.15*(cols-1))/cols
    card_h = 1.7 if rows<=2 else 1.45
    y_start = 1.3 if rows<=2 else 1.25
    y_step = card_h + 0.25 if rows<=2 else card_h + 0.2
    for i,(sem,scn,sen) in enumerate(subtypes):
        col=i%cols;row=i//cols
        x=0.3+col*(card_w+0.15);y=y_start+row*y_step
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(card_w),Inches(card_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
        tb(s,x+0.1,y+0.1,card_w-0.2,0.55,sem,sz=30,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.75,card_w-0.2,0.4,scn,sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+1.15,card_w-0.2,0.3,sen,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    return s

def word_card_read(w,py,en,sent,img,color=MAGENTA):
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=color,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=color)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,img,color=MAGENTA):
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=color)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=color)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    return s

def master_intro_slide(name_cn, name_en, country_cn, country_en, years, summary_cn, summary_en, color, img_lb):
    s=ns();bg(s,CREAM);hb(s,f"🎨 {name_cn}  {name_en}",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.4),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(3)
    ib(s,0.6,1.15,4.0,3.8,img_lb)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(0.95),Inches(4.7),Inches(4.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=color;sh2.line.fill.background()
    tb(s,5.2,1.1,4.4,0.6,name_cn,sz=30,b=True,c=WHITE)
    tb(s,5.2,1.7,4.4,0.4,name_en,sz=15,c=WARM)
    tb(s,5.2,2.25,4.4,0.4,f"🌍 国家 / Country",sz=13,b=True,c=YELLOW)
    tb(s,5.2,2.6,4.4,0.4,f"{country_cn}  {country_en}",sz=18,b=True,c=WHITE)
    tb(s,5.2,3.1,4.4,0.4,"🗓️ 年代 / Years",sz=13,b=True,c=YELLOW)
    tb(s,5.2,3.45,4.4,0.4,years,sz=18,b=True,c=WHITE)
    tb(s,5.2,3.95,4.4,0.4,"💡 一句话 / In one line",sz=13,b=True,c=YELLOW)
    tb(s,5.2,4.3,4.4,0.4,summary_cn,sz=15,b=True,c=WHITE)
    tb(s,5.2,4.7,4.4,0.35,summary_en,sz=11,c=WARM)
    return s

n=0

# 1 COVER
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.7,"Follow the Masters",sz=32,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.45,"跟着大师学画画",sz=22,c=MAGENTA,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.55),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=MAGENTA;sh.line.color.rgb=YELLOW;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.85),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=YELLOW;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.0,2.8,0.4,"DAY 2",sz=16,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"🖼️",sz=52,a=PP_ALIGN.CENTER)
ap(tf,"跟着大师学画画",sz=18,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"FOLLOW THE MASTERS",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
for x,y,r,c in [(1.2,1.9,0.6,PICASSO),(1.5,3.8,0.5,MATISSE),(7.8,1.6,0.55,VANGOGH),(8.2,3.6,0.5,SKY),(1.0,4.7,0.35,GREEN_L),(8.6,4.6,0.35,PURPLE)]:
    dot(s,x,y,r,c)
tb(s,1,5.05,8,0.4,"🎨 毕加索 · 马蒂斯 · 梵高  Picasso · Matisse · Van Gogh",sz=14,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
pn(s,n)

# 2 SCHEDULE
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","认识三位大师 + 风格 + 代表作",MAGENTA),
    ("Session 2  下午","2:00-2:45","复习 + 我会认 4 词 / 我会写 3 词",SKY),
    ("Session 3  下午","3:00-4:30","动手做大师作品 + 画展",CORAL)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=WARM)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=WHITE)
pn(s,n)

# 3 OBJECTIVES
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"🎨 内容目标  Content:",sz=20,b=True,c=MAGENTA)
tf=tb(s,0.7,1.4,9,1.3,"1. 认识三位大师：毕加索、马蒂斯、梵高",sz=15,c=DARK)
ap(tf,"2. 欣赏他们的代表作品",sz=15,c=DARK)
ap(tf,"3. 发现他们不同的艺术风格",sz=15,c=DARK)
ap(tf,"4. 学习用大师的方法表达自己",sz=15,c=DARK)
tb(s,0.5,3.0,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=SKY)
tb(s,0.7,3.5,9,0.4,"👀 我会认：画家  颜色  形状  艺术",sz=15,b=True,c=DARK)
tb(s,0.7,4.0,9,0.4,"✍️ 我会写：画家  大师  艺术",sz=15,b=True,c=DARK)
tb(s,0.5,4.65,9,0.5,"🖌️ 实践目标：完成 D2 booklet + 1 幅作品",sz=15,c=CORAL)
pn(s,n)

# 4 SESSION 1 DIVIDER
div("Session 1  上午","认识三位大师\n🎨 毕加索  ·  ✂️ 马蒂斯  ·  🌻 梵高",MAGENTA,"🖼️");n+=1

# 5 What is a Master
s=ns();n+=1;bg(s,CREAM);hb(s,"💡 什么是大师？  What is a Master?",MAGENTA)
# Slim definition callout
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(0.95),Inches(8.4),Inches(0.75))
sh.fill.solid();sh.fill.fore_color.rgb=MAGENTA;sh.line.fill.background()
tb(s,1.0,0.98,8.0,0.45,"大师 = 有自己的风格的艺术家",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.0,1.4,8.0,0.3,"A Master = an artist with their own STYLE",sz=12,c=YELLOW,a=PP_ALIGN.CENTER)
# 3 example cards: famous painting + signature feature label
examples=[
    ("📷 《哭泣的女人》Weeping Woman","《哭泣的女人》","变形的脸","Twisted face","= 毕加索",PICASSO,"en.wikipedia.org/wiki/The_Weeping_Woman","weeping_woman"),
    ("📷 《伊卡洛斯》Icarus (黑色剪影 + 红心)","《伊卡洛斯》","彩色剪纸","Cut-out shapes","= 马蒂斯",MATISSE,"搜「Matisse Icarus Jazz」","icarus"),
    ("📷 《星夜》The Starry Night (旋转的天空)","《星夜》","旋转线条","Swirly lines","= 梵高",VANGOGH,"en.wikipedia.org/wiki/The_Starry_Night","starry_night"),
]
for i,(img_lb,title,feat_cn,feat_en,who,cl,url,key) in enumerate(examples):
    x=0.3+i*3.2
    # card background
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(3.0),Inches(3.35))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    # image placeholder (top)
    ib(s,x+0.15,1.95,2.7,1.95,img_lb,url=url,key=key)
    # painting title strip
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.15),Inches(3.95),Inches(2.7),Inches(0.35))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+0.2,3.97,2.6,0.32,title,sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # feature → master
    tb(s,x+0.15,4.4,2.7,0.4,f"看! {feat_cn}",sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,4.78,2.7,0.3,who,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# 6 3-master overview
s=ns();n+=1;bg(s,CREAM);hb(s,"👥 今天的三位大师  Today's 3 Masters",MAGENTA)
tb(s,0.4,0.85,9,0.35,"看一看 — 你认识他们吗？",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
masters=[
    ("🎭","毕加索","Picasso","西班牙 Spain","1881-1973",PICASSO),
    ("✂️","马蒂斯","Matisse","法国 France","1869-1954",MATISSE),
    ("🌻","梵高","Van Gogh","荷兰 Netherlands","1853-1890",VANGOGH),
]
for i,(em,cn,en,country,years,cl) in enumerate(masters):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.3),Inches(3.0),Inches(3.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    ib(s,x+0.2,1.45,2.6,1.7,f"📷 {cn} 头像")
    tb(s,x+0.1,3.2,2.8,0.5,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.65,2.8,0.45,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.1,2.8,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.4,2.8,0.3,country,sz=11,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.7,2.8,0.3,years,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# 7 Picasso intro
s=master_intro_slide("毕加索","Pablo Picasso","西班牙","Spain","1881-1973",
    "他改变了我们看人脸的方式。",
    "He changed the way we see faces.",
    PICASSO,"📷 毕加索头像 Picasso portrait");n+=1;pn(s,n)

# 8 Picasso style — image example on left, concept cards on right
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 毕加索的风格 — 变形的人脸",PICASSO)
tb(s,0.4,0.9,9.2,0.35,"Cubism  立体派 — 把不同角度放在一起！",sz=15,b=True,c=PICASSO,a=PP_ALIGN.CENTER)
# LEFT: big example painting + caption strip
ib(s,0.3,1.35,4.4,2.6,"📷 《多拉·玛尔的肖像》Portrait of Dora Maar",url="en.wikipedia.org/wiki/Portrait_of_Dora_Maar",key="dora_maar")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.0),Inches(4.4),Inches(0.42))
sh.fill.solid();sh.fill.fore_color.rgb=PICASSO;sh.line.fill.background()
tb(s,0.4,4.04,4.2,0.35,"《多拉·玛尔的肖像》· 1937",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# RIGHT: 4 concept cards 2x2 — anchor each to the painting
icons=[("👁️","正面的眼睛","Front eye"),("👃","侧面的鼻子","Side nose"),("👄","变形的嘴","Twisted mouth"),("🟦","几何形状","Geometric shapes")]
for i,(em,cn,en) in enumerate(icons):
    col=i%2;row=i//2
    x=4.85+col*2.45;y=1.35+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.35),Inches(1.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PICASSO;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.08,2.15,0.5,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.62,2.15,0.4,cn,sz=14,b=True,c=PICASSO,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.02,2.15,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Bottom slim inquiry strip
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.6),Inches(9.4),Inches(0.55))
sh.fill.solid();sh.fill.fore_color.rgb=PICASSO;sh.line.fill.background()
tb(s,0.4,4.65,9.2,0.5,"🤔 看这幅画 — 你能找出几个「不一样的角度」？  Find the angles!",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)
notes(s,"老师备课:\n• 立体派 (Cubism) 的核心: 同一张脸里有「正面」+「侧面」混在一起\n• 传统画法只能选一个角度, 毕加索打破了这个规则\n• 找角度提示: 眼睛(正面) / 鼻子(侧面) / 嘴(扭曲) / 脸(几何块)\n• 追问: 「平常我们只能看到正面 OR 侧面, 毕加索同时画两个 — 你试过吗？」\n• 趣闻: 毕加索全名超长 — Pablo Diego José Francisco de Paula Juan Nepomuceno María de los Remedios Cipriano de la Santísima Trinidad Ruiz Picasso!")

# 9a Observe Three Musicians (before teacher explains)
s=observe_slide("🎭","三个音乐家","Three Musicians",PICASSO,
    "📷 三个音乐家 Three Musicians",
    ["你看到几个人？他们在做什么？",
     "你能找到几个乐器？分别是什么？",
     "形状像什么？(三角形 / 方形 / 圆?)",
     "桌下面藏着一个小动物 — 你能找到吗？"],
    url="en.wikipedia.org/wiki/Three_Musicians",key="three_musicians")
n+=1;pn(s,n)

# 9 Picasso work 1
s=example_slide("🎭","毕加索","三个音乐家","Three Musicians",
    "Pablo Picasso · 1921 年 · 综合立体派",
    "三个音乐家 — 你能找到几个乐器？\n3 musicians — can you find their instruments?",
    "📷 三个音乐家 Three Musicians",PICASSO,
    url="en.wikipedia.org/wiki/Three_Musicians",key="three_musicians");n+=1;pn(s,n)
notes(s,"老师备课:\n• 1921 年作品, 立体派的彩色、欢乐版本 (「综合立体派」)\n• 三个角色: 左边 Pierrot (戴白帽小丑) + 中间 Harlequin (花格小丑) + 右边修道士\n• 像剪贴画 — 几何形状拼起来\n• 找一找: 三个乐器 (单簧管 / 吉他 / 乐谱), 桌下的小狗也是几何形\n• 趣闻: 这幅画有两个版本 — 一个在 MoMA, 一个在费城\n• 追问: 「他们在演奏什么？」「形状像什么？」「能不能用纸剪出这样的人？」")

# 10a Observe Weeping Woman
s=observe_slide("😢","哭泣的女人","Weeping Woman",PICASSO,
    "📷 哭泣的女人 Weeping Woman",
    ["这个人在做什么？你怎么知道？",
     "她的眼睛是什么形状？(像不像眼泪？)",
     "嘴是什么形状？锯齿让你想到什么？",
     "主要用了什么颜色？这些颜色让你感觉?"],
    url="en.wikipedia.org/wiki/The_Weeping_Woman",key="weeping_woman")
n+=1;pn(s,n)

# 10 Picasso work 2
s=example_slide("🎭","毕加索","哭泣的女人","Weeping Woman",
    "Pablo Picasso · 1937 年 · 蓝色和绿色",
    "你觉得她伤心吗？为什么？\nDo you think she is sad? Why?",
    "📷 哭泣的女人 Weeping Woman",PICASSO,
    url="en.wikipedia.org/wiki/The_Weeping_Woman",key="weeping_woman");n+=1;pn(s,n)
notes(s,"老师备课:\n• 1937 年, 画家女友 Dora Maar 的肖像\n• 同一年毕加索画了著名的《格尔尼卡》— 反战, 因为纳粹刚轰炸西班牙小镇 Guernica, 死了 1000 多平民\n• 锯齿形眼泪 + 蓝绿色调 = 痛苦 + 哭泣\n• 嘴巴是锯齿状, 眼睛形状像眼泪\n• 追问: 「眼泪在哪里？锯齿让你想到什么？」\n• 链接 D1: 蓝/绿 = 难过的颜色")

# 11 Picasso try-it — concrete 5-step recipe
s=ns();n+=1;bg(s,CREAM);hb(s,"✏️ 毕加索 试一试  Try Picasso!",PICASSO)
# Title strip
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(9.4),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=PICASSO;sh.line.fill.background()
tb(s,0.4,1.0,7.0,0.4,"⏱️ 5 分钟画一张毕加索人脸！",sz=17,b=True,c=WHITE)
tb(s,7.0,1.05,2.6,0.35,"5-min Picasso face",sz=11,c=YELLOW,a=PP_ALIGN.RIGHT)
# LEFT: 5-step recipe
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.6),Inches(5.3),Inches(3.55))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PICASSO;sh.line.width=Pt(2.5)
tb(s,0.5,1.7,5.0,0.4,"👉 跟着做  5 Steps",sz=15,b=True,c=PICASSO)
steps=[
    ("1️⃣","画一个圆 (头)","Draw a circle (head)"),
    ("2️⃣","中间画一条线，分两半","Line down the middle"),
    ("3️⃣","左边: 画正面的圆眼睛 + 嘴","Left half: front eye + mouth"),
    ("4️⃣","右边: 三角形鼻子 + 半张嘴","Right half: triangle nose + half mouth"),
    ("5️⃣","填 2-3 种颜色","Fill with 2-3 colors"),
]
for i,(num,cn,en) in enumerate(steps):
    y=2.15+i*0.58
    tb(s,0.55,y,0.5,0.4,num,sz=18,b=True,c=PICASSO)
    tb(s,1.1,y+0.02,4.4,0.32,cn,sz=13,b=True,c=DARK)
    tb(s,1.1,y+0.32,4.4,0.25,en,sz=10,c=GRAY)
# RIGHT: demo image
ib(s,5.8,1.6,3.9,3.55,"📷 老师示范图 / Demo:\n圆 + 中线 + 一边正面 + 一边侧面")
pn(s,n)

# 12 Matisse intro
s=master_intro_slide("马蒂斯","Henri Matisse","法国","France","1869-1954",
    "他用剪刀「画画」。",
    "He painted with scissors.",
    MATISSE,"📷 马蒂斯头像 Matisse portrait");n+=1;pn(s,n)
notes(s,"老师备课:\n• 法国画家, 与毕加索齐名的 20 世纪大师\n• 他追求「颜色的自由运用」— 那年代画家用色还受真实限制 (天必须蓝, 草必须绿)\n• 比印象派和梵高更进一步, 敢用「不真实」的颜色\n• 老了 + 生病后开始用剪刀做剪纸, 这是他第二次艺术革命\n• 追问: 「你有没有用过「不像真的」的颜色画过画？」")

# 13 Matisse style — image example on left, story + shapes on right
s=ns();n+=1;bg(s,CREAM);hb(s,"✂️ 马蒂斯的风格 — 彩色剪纸",MATISSE)
tb(s,0.4,0.9,9.2,0.35,"Cut-Outs 剪纸艺术 — 用剪刀代替画笔！",sz=15,b=True,c=MATISSE,a=PP_ALIGN.CENTER)
# LEFT: big example painting
ib(s,0.3,1.35,4.4,2.6,"📷 《波利尼西亚，海》Polynesia, The Sea (海洋生物剪纸)",url="搜「Matisse Polynesia the Sea」",key="polynesia_sea")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.0),Inches(4.4),Inches(0.42))
sh.fill.solid();sh.fill.fore_color.rgb=MATISSE;sh.line.fill.background()
tb(s,0.4,4.04,4.2,0.35,"《波利尼西亚，海》· 1946",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# RIGHT TOP: story panel
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.35),Inches(4.85),Inches(1.4))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=MATISSE;sh.line.width=Pt(2)
tb(s,5.0,1.42,4.6,0.4,"🌟 他的故事 His Story",sz=14,b=True,c=MATISSE)
tf=tb(s,5.0,1.82,4.6,0.4,"老了，不能站着画画 →",sz=13,c=DARK)
ap(tf,"拿起剪刀剪形状！",sz=13,b=True,c=MATISSE)
ap(tf,"Old + sick → scissors + paper",sz=10,c=GRAY)
# RIGHT BOTTOM: 4 shape cards 2x2
shapes=[("🟦","蓝色方块",SKY),("⭐","黄色星星",YELLOW),("🌿","绿色叶子",GREEN_L),("💗","粉色心",MAGENTA)]
for i,(em,cn,cl) in enumerate(shapes):
    col=i%2;row=i//2
    x=4.85+col*2.45;y=2.85+row*0.7
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.35),Inches(0.6))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,x+0.1,y+0.1,0.55,0.4,em,sz=18,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.7,y+0.13,1.6,0.4,cn,sz=13,b=True,c=WHITE)
# Bottom slim inquiry strip
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(4.3),Inches(4.85),Inches(0.55))
sh.fill.solid();sh.fill.fore_color.rgb=MATISSE;sh.line.fill.background()
tb(s,5.0,4.35,4.65,0.5,"🤔 你能剪出哪些形状？  What shapes?",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)
notes(s,"老师备课:\n• 马蒂斯把画面简化到像儿童画 — 没有光影、没有立体感, 只有形状 + 颜色\n• 互补色技巧 (色环): 红+绿 / 黄+紫 / 蓝+橙 — 放一起视觉冲击强烈\n• 趣闻: 他的名作《红房间》原本是蓝色, 后来才改成红色 — 因为红色和窗外绿色互补, 更有冲击力\n• 《波利尼西亚, 海》灵感来自他在大溪地的旅行, 看到的鱼和珊瑚\n• 追问: 「你能用 3 种颜色, 表达「夏天」吗？」")

# 14a Observe Snail
s=observe_slide("🐌","蜗牛","The Snail",MATISSE,
    "📷 蜗牛 The Snail (彩色方块螺旋)",
    ["这是什么？(让你猜!)",
     "颜色按什么顺序排？(红→橙→黄→…?)",
     "形状是规则的, 还是随意的？",
     "哪里是「空白」？为什么留白？"],
    url="en.wikipedia.org/wiki/The_Snail",key="snail")
n+=1;pn(s,n)

# 14 Matisse work 1
s=example_slide("✂️","马蒂斯","蜗牛","The Snail",
    "Henri Matisse · 1953 年 · 剪纸作品",
    "彩色方块拼成一只蜗牛！\nA snail made of colored squares!",
    "📷 蜗牛 The Snail (彩色方块螺旋)",MATISSE,
    url="en.wikipedia.org/wiki/The_Snail",key="snail");n+=1;pn(s,n)
notes(s,"老师备课:\n• 1953 年作品, 马蒂斯生命最后阶段, 已经躺在床上做艺术 (84 岁)\n• 形状像蜗牛壳的螺旋: 红 → 橙 → 黄 → 绿 → 蓝 → 紫\n• 即使是简单彩色方块也是艺术 — 重点在颜色 + 形状的关系\n• 没有线条, 没有透视 — 但很有动感\n• 追问: 「你能看出蜗牛吗？」「为什么不画一只真的蜗牛？」")

# 15a Observe Icarus
s=observe_slide("⭐","伊卡洛斯","Icarus",MATISSE,
    "📷 伊卡洛斯 Icarus",
    ["黑色的形状是什么？",
     "中间那颗红色的圆 — 你觉得是什么？",
     "你看到几颗黄色的星？",
     "蓝色背景代表什么？(天空? 海? 夜晚?)"],
    url="搜「Matisse Icarus Jazz」",key="icarus")
n+=1;pn(s,n)

# 15 Matisse work 2
s=example_slide("✂️","马蒂斯","伊卡洛斯","Icarus",
    "Henri Matisse · 1944 年 · 剪纸 + 油彩",
    "身体只是一个剪影 — 心是红色的星！\nBody = silhouette · heart = red star!",
    "📷 伊卡洛斯 Icarus",MATISSE,
    url="搜「Matisse Icarus Jazz」",key="icarus");n+=1;pn(s,n)
notes(s,"老师备课:\n• 1944 年作品, 来自他的「Jazz」剪纸书\n• 故事: 希腊神话 — 伊卡洛斯用蜡做的翅膀飞, 太靠近太阳, 蜡融化, 他坠落\n• 黑色 = 身体剪影 / 红色 = 心 / 黄色点 = 星星 / 蓝色 = 天空\n• 不强调坠落悲剧, 重点是「剪影 + 一个亮点」的画法\n• 追问: 「红色代表什么？为什么心是亮的？」(链接 D1 心情色)")

# 16 Matisse try-it — concrete 5-step recipe
s=ns();n+=1;bg(s,CREAM);hb(s,"✂️ 马蒂斯 试一试  Try Matisse!",MATISSE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(9.4),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=MATISSE;sh.line.fill.background()
tb(s,0.4,1.0,7.0,0.4,"⏱️ 5 分钟做一幅马蒂斯小剪纸！",sz=17,b=True,c=WHITE)
tb(s,7.0,1.05,2.6,0.35,"5-min Matisse cut-out",sz=11,c=YELLOW,a=PP_ALIGN.RIGHT)
# LEFT: 5-step recipe
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.6),Inches(5.3),Inches(3.55))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MATISSE;sh.line.width=Pt(2.5)
tb(s,0.5,1.7,5.0,0.4,"👉 跟着做  5 Steps",sz=15,b=True,c=MATISSE)
steps=[
    ("1️⃣","选 3 种颜色的彩纸","Pick 3 paper colors"),
    ("2️⃣","不画线，直接拿剪刀剪","No drawing — just cut"),
    ("3️⃣","剪 3 个形状: 圆 / 星 / 心","Cut 3 shapes: circle/star/heart"),
    ("4️⃣","在白纸上摆好位置","Arrange on white paper"),
    ("5️⃣","喜欢就用胶水贴上！","Glue down what you like!"),
]
for i,(num,cn,en) in enumerate(steps):
    y=2.15+i*0.58
    tb(s,0.55,y,0.5,0.4,num,sz=18,b=True,c=MATISSE)
    tb(s,1.1,y+0.02,4.4,0.32,cn,sz=13,b=True,c=DARK)
    tb(s,1.1,y+0.32,4.4,0.25,en,sz=10,c=GRAY)
# RIGHT: demo image
ib(s,5.8,1.6,3.9,3.55,"📷 老师示范图 / Demo:\n彩纸 + 剪刀 + 简单形状拼贴")
pn(s,n)

# 17 Van Gogh intro
s=master_intro_slide("梵高","Vincent van Gogh","荷兰","Netherlands","1853-1890",
    "他的颜色会跳舞。",
    "His colors dance.",
    VANGOGH,"📷 梵高头像 Van Gogh portrait");n+=1;pn(s,n)
notes(s,"老师备课:\n• 荷兰画家, 一生只活了 37 岁\n• 27 岁才开始画画, 10 年内画了 2000 多幅作品 (画 / 素描)\n• 生前只卖出 1 幅画! 死后才被全世界认可\n• 他不是「疯子」, 是用颜色 + 线条表达内心强烈情绪\n• 比印象派更进一步 — 不只画看到的, 更要画感受到的\n• 追问: 「你今天的心情是什么颜色？」")

# 18 Van Gogh style — image example on left, concept cards on right
s=ns();n+=1;bg(s,CREAM);hb(s,"🌻 梵高的风格 — 有情绪的颜色和线条",VANGOGH)
tb(s,0.4,0.9,9.2,0.35,"短粗笔触 + 旋转线条 = 情绪",sz=15,b=True,c=VANGOGH,a=PP_ALIGN.CENTER)
# LEFT: big example painting
ib(s,0.3,1.35,4.4,2.6,"📷 《罗讷河上的星夜》Starry Night Over the Rhône (旋转笔触)",url="en.wikipedia.org/wiki/Starry_Night_Over_the_Rhône",key="starry_rhone")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.0),Inches(4.4),Inches(0.42))
sh.fill.solid();sh.fill.fore_color.rgb=VANGOGH;sh.line.fill.background()
tb(s,0.4,4.04,4.2,0.35,"《罗讷河上的星夜》· 1888",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
# RIGHT: 4 concept cards 2x2 — anchor each to the painting
icons=[("〰️","旋转线条","Swirly lines"),("🖌️","短粗笔触","Short strokes"),("🌈","明亮颜色","Bright colors"),("❤️","表达心情","Show feelings")]
for i,(em,cn,en) in enumerate(icons):
    col=i%2;row=i//2
    x=4.85+col*2.45;y=1.35+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.35),Inches(1.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=VANGOGH;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.08,2.15,0.5,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.62,2.15,0.4,cn,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.02,2.15,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Bottom slim inquiry strip
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.6),Inches(9.4),Inches(0.55))
sh.fill.solid();sh.fill.fore_color.rgb=VANGOGH;sh.line.fill.background()
tb(s,0.4,4.65,9.2,0.5,"🤔 看线条 — 你觉得他开心还是难过？  Happy or sad?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)
notes(s,"老师备课:\n• 1888 年画的, 梵高在阿尔勒 (Arles) 河边写生\n• 与 1889 年的《星夜》是「姐妹作」\n• 黄色光点 = 天上的星 + 河边煤气灯的倒影\n• 短粗笔触表现夜晚的空气在「呼吸」, 风像在动\n• 期待: 学生从颜色亮 / 暗 + 线条快 / 慢, 来判断心情\n• 关键: 梵高一生很苦, 但用亮色画 — 颜色不是直接 = 心情, 是表达\n• 追问: 「天上的星星和河里的灯影分得清吗？」「为什么画家用旋转的笔？」")

# 19a Observe Sunflowers
s=observe_slide("🌻","向日葵","Sunflowers",VANGOGH,
    "📷 向日葵 Sunflowers",
    ["数一数 — 一共几朵向日葵？",
     "花都长得一样吗？(有的开了, 有的谢了?)",
     "主要用了什么颜色？",
     "花瓶上写了什么？(有「Vincent」签名!)"],
    url="en.wikipedia.org/wiki/Sunflowers_(Van_Gogh_series)",key="sunflowers")
n+=1;pn(s,n)

# 19 Van Gogh work 1
s=example_slide("🌻","梵高","向日葵","Sunflowers",
    "Vincent van Gogh · 1888 年 · 油画",
    "黄色 + 黄色 — 像太阳！\nYellow + yellow — like the sun!",
    "📷 向日葵 Sunflowers",VANGOGH,
    url="en.wikipedia.org/wiki/Sunflowers_(Van_Gogh_series)",key="sunflowers");n+=1;pn(s,n)
notes(s,"老师备课:\n• 1888 年画的, 准备给好朋友 Gauguin (高更) 看 — 高更要来阿尔勒和梵高同住\n• 梵高一系列向日葵共 11 幅, 这是最有名的几幅之一\n• 黄色 = 友谊 + 希望 (梵高的色彩象征)\n• 你看到的花有的开了, 有的谢了 — 像生命过程\n• 追问: 「你看到几朵？它们都一样吗？」「你最喜欢哪一朵？」")

# 20a Observe Starry Night
s=observe_slide("🌌","星夜","The Starry Night",VANGOGH,
    "📷 星夜 The Starry Night",
    ["是白天还是晚上？怎么知道？",
     "数一数 — 多少颗星星？多少个月亮？",
     "天空在动吗？画家怎么让它「动」起来?",
     "下面地上有什么？(村庄? 教堂? 树?)"],
    url="en.wikipedia.org/wiki/The_Starry_Night",key="starry_night")
n+=1;pn(s,n)

# 20 Van Gogh work 2
s=example_slide("🌻","梵高","星夜","The Starry Night",
    "Vincent van Gogh · 1889 年 · 油画",
    "天上的星星在转 — 风在动！\nStars are spinning · wind is moving!",
    "📷 星夜 The Starry Night",VANGOGH,
    url="en.wikipedia.org/wiki/The_Starry_Night",key="starry_night");n+=1;pn(s,n)
notes(s,"老师备课:\n• 1889 年在精神病院 (Saint-Rémy) 画的, 梵高从窗户看出去\n• 这不是真实的景, 是想象 + 记忆 — 中间的尖塔是村里的教堂 (荷兰记忆), 但景是法国南部的山\n• 画里有 11 颗星 + 1 个月亮, 螺旋的天空像漩涡\n• 趣闻: 现代天文学家发现, 螺旋的形状非常像真实的星系！梵高在 1889 年就「画对了」\n• 这是世界上最有名的画之一, 可以问「你在哪里看过？」(海报、马克杯、T 恤……)\n• 追问: 「天上的星星会动吗？为什么画家画成这样？」「你做梦时, 看到的世界和真实一样吗？」")

# 21 Van Gogh try-it — concrete 5-step recipe
s=ns();n+=1;bg(s,CREAM);hb(s,"🌻 梵高 试一试  Try Van Gogh!",VANGOGH)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(9.4),Inches(0.5))
sh.fill.solid();sh.fill.fore_color.rgb=VANGOGH;sh.line.fill.background()
tb(s,0.4,1.0,7.0,0.4,"⏱️ 5 分钟画一张「心情线」！",sz=17,b=True,c=DARK)
tb(s,7.0,1.05,2.6,0.35,"5-min mood lines",sz=11,c=DARK,a=PP_ALIGN.RIGHT)
# LEFT: 5-step recipe
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.6),Inches(5.3),Inches(3.55))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=VANGOGH;sh.line.width=Pt(2.5)
tb(s,0.5,1.7,5.0,0.4,"👉 跟着做  5 Steps",sz=15,b=True,c=DARK)
steps=[
    ("1️⃣","选 1 个心情: 开心 / 生气 / 难过","Pick a mood"),
    ("2️⃣","选 2-3 种颜色 (黄=开心 红=生气 蓝=难过)","Pick 2-3 colors for that mood"),
    ("3️⃣","不画具体东西，只画线条!","No objects — just lines!"),
    ("4️⃣","短而粗的笔触，像雨点","Short bold strokes, like raindrops"),
    ("5️⃣","让线条旋转 / 跳舞 — 心情就来了!","Swirl them — feeling appears!"),
]
for i,(num,cn,en) in enumerate(steps):
    y=2.15+i*0.58
    tb(s,0.55,y,0.5,0.4,num,sz=18,b=True,c=VANGOGH)
    tb(s,1.1,y+0.02,4.4,0.32,cn,sz=12,b=True,c=DARK)
    tb(s,1.1,y+0.32,4.4,0.25,en,sz=10,c=GRAY)
# RIGHT: demo image
ib(s,5.8,1.6,3.9,3.55,"📷 老师示范图 / Demo:\n短粗笔触 + 旋转线条 = 心情")
pn(s,n)

# 22 Comparison table
s=ns();n+=1;bg(s,CREAM);hb(s,"📊 三位大师对比  3 Masters Comparison",MAGENTA)
tb(s,0.4,0.85,9,0.35,"看表格 — 找一找他们的不同！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(4,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.7));t=ts.table
t.columns[0].width=Inches(1.6)
t.columns[1].width=Inches(2.6)
t.columns[2].width=Inches(2.6)
t.columns[3].width=Inches(2.6)
rows=[["大师 Master","🎭 毕加索 Picasso","✂️ 马蒂斯 Matisse","🌻 梵高 Van Gogh"],
      ["🌍 国家 Country","西班牙 Spain","法国 France","荷兰 Netherlands"],
      ["🎨 风格 Style","立体派 / 变形人脸","彩色剪纸 / 形状","旋转线条 / 情绪色彩"],
      ["🖼️ 代表作 Famous","《三个音乐家》\n《哭泣的女人》","《蜗牛》\n《伊卡洛斯》","《向日葵》\n《星夜》"]]
header_colors=[MAGENTA,PICASSO,MATISSE,VANGOGH]
for r,rd in enumerate(rows):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 12);rn.font.bold=(r==0 or c==0)
        if r==0:
            rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=header_colors[c]
        elif c==0:
            rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        else:
            rn.font.color.rgb=DARK
            if r%2==1:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF7,0xF3,0xE9)
pn(s,n)

# 23 Whose style guess
s=ns();n+=1;bg(s,CREAM);hb(s,"🕵️ 这是谁的风格？  Whose Style?",MAGENTA)
tb(s,0.4,0.85,9,0.35,"看图，猜大师！(口头) Look and guess the master!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
mysteries=[("？1","旋转的星空","Swirly stars",VANGOGH),("？2","三角形的脸","Triangle face",PICASSO),("？3","彩色剪纸","Cut-out shapes",MATISSE)]
for i,(num,hint,en,cl) in enumerate(mysteries):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.3),Inches(3.0),Inches(3.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.4,2.8,0.5,num,sz=26,b=True,c=cl,a=PP_ALIGN.CENTER)
    ib(s,x+0.2,1.95,2.6,2.0,"📷 神秘画作 Mystery painting")
    tb(s,x+0.1,4.05,2.8,0.4,f"线索: {hint}",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.45,2.8,0.3,f"Clue: {en}",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.8,0.3,"是谁？Who is it?",sz=12,b=True,c=cl,a=PP_ALIGN.CENTER)
pn(s,n)

# 24 SESSION 2 DIVIDER
div("Session 2  下午","复习 + 语言目标\n👀 我会认 4 词  ·  ✍️ 我会写 3 词",SKY,"📖");n+=1

# 25 Quick review match
s=ns();n+=1;bg(s,CREAM);hb(s,"🔄 快速复习  Quick Review — 连大师 ↔ 风格",SKY)
tb(s,0.4,0.85,9,0.35,"把大师和风格连起来 (口头)  Match the master to their style!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
m_left=[("🎭 毕加索",PICASSO),("✂️ 马蒂斯",MATISSE),("🌻 梵高",VANGOGH)]
m_right=["旋转线条 + 明亮颜色","立体派 / 变形人脸","彩色剪纸 / 形状"]
for i,(em,cl) in enumerate(m_left):
    y=1.3+i*1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(y),Inches(3.4),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.95,y+0.2,3.2,0.5,em,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
for i,w in enumerate(m_right):
    y=1.3+i*1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.8),Inches(y),Inches(3.4),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SKY;sh.line.width=Pt(2)
    tb(s,5.95,y+0.25,3.2,0.5,w,sz=16,b=True,c=SKY,a=PP_ALIGN.CENTER)
tb(s,4.4,2.7,1.2,0.6,"?",sz=44,b=True,c=CORAL,a=PP_ALIGN.CENTER)
tb(s,0.4,4.85,9,0.4,"举手回答 — 谁的风格是哪个？",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# 26-29 我会认 word cards
read_words=[
    ("画家","huà jiā","painter","毕加索是一个画家。","📷 画家在画画",MAGENTA),
    ("颜色","yán sè","color","梵高用很多颜色画画。","📷 各种颜色",VANGOGH),
    ("形状","xíng zhuàng","shape","马蒂斯剪了很多形状。","📷 各种形状",MATISSE),
    ("艺术","yì shù","art","我喜欢艺术。","📷 艺术作品",PICASSO),
]
for w,py,en,sent,img,cl in read_words:
    s=word_card_read(w,py,en,sent,img,cl);n+=1;pn(s,n)

# 30 Word games
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",SKY)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","老师说词\n学生拍正确字卡",WARM),
    ("2️⃣","配大师\nMaster Match","看大师图\n说对应词语",RGBColor(0xE3,0xF2,0xFD)),
    ("3️⃣","抢字卡\nGrab the Card","老师读词\n第一个抢到赢",RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣","我演你猜\nActing Game","做动作 / 表情\n大家猜词",RGBColor(0xFC,0xE4,0xEC)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# 31-33 我会写 word cards
write_words=[
    ("画家","huà jiā","painter","📷 画家",MAGENTA),
    ("大师","dà shī","master","📷 大师",PURPLE),
    ("艺术","yì shù","art","📷 艺术",PICASSO),
]
for w,py,en,img,cl in write_words:
    s=word_card_write(w,py,en,img,cl);n+=1;pn(s,n)

# 34 SESSION 3 DIVIDER
div("Session 3  下午","动手做大师作品！  Make Masters' Art!\n🎨 拼贴  ·  🌻 向日葵  ·  🎭 立体自画像",CORAL,"🖌️");n+=1

# 35 Booklet
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成 D2 练习册  Day 2 Booklet',CORAL)
ib(s,0.4,0.9,9.2,4.3,"📷 D2 练习册截图 / Booklet pages")
pn(s,n)

# 36 Projects overview — 2 levels
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 今天做什么？  Two Project Levels",CORAL)
tb(s,0.4,0.9,9,0.4,"老师会根据你的年级选一个！Teacher will pick for your level.",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
projects=[
    ("低 Level\n(K-2)","✂️ 马蒂斯拼贴 OR 🌻 梵高向日葵","Matisse cut-out OR Sunflowers","老师二选一 — 看孩子兴趣\nTeacher picks based on class",MATISSE,VANGOGH),
    ("高 Level\n(3-5)","🎭 毕加索立体自画像","Picasso Cubist Self-Portrait","画 + 写一句介绍\nDraw + write 1 sentence",PICASSO,PICASSO),
]
for i,(lvl,nm,en,d,ltcl,cl) in enumerate(projects):
    x=0.5+i*4.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(4.3),Inches(3.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(4)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.25),Inches(1.55),Inches(1.4),Inches(0.7))
    sh2.fill.solid();sh2.fill.fore_color.rgb=ltcl;sh2.line.fill.background()
    tf=tb(s,x+0.3,1.62,1.3,0.5,lvl.split('\n')[0],sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    ap(tf,lvl.split('\n')[1],sz=10,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,2.35,4.0,0.6,nm,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,2.95,4.0,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.5,3.35,3.4,1.2,"📷 作品示范")
    ls=d.split('\n')
    tf=tb(s,x+0.2,4.65,4.0,0.25,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# 37 Project 低 Level — Matisse cut-out
s=ns();n+=1;bg(s,CREAM);hb(s,"✂️ 低 Level: 马蒂斯彩纸拼贴  Matisse Cut-Out",MATISSE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=MATISSE;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"🎨 彩色纸 (3-4 色)  Colored paper",sz=13,c=DARK)
ap(tf,"✂️ 剪刀  Scissors",sz=13,c=DARK)
ap(tf,"🧴 胶水 / 胶棒  Glue stick",sz=13,c=DARK)
ap(tf,"🖼️ 画板 / 大白纸  Board / paper",sz=13,c=DARK)
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=YELLOW;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"🗣️ 句型 Sentence Frames",sz=14,b=True,c=DARK)
tf_k=tb(s,0.4,3.85,4.4,1.4,"· 我用了「红色 / 蓝色 / 黄色」。",sz=13,c=DARK)
ap(tf_k,"· 我剪了一个「圆 / 三角形」。",sz=13,c=DARK)
ap(tf_k,"· 这是我的拼贴画！",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SKY;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 选 3-4 种颜色的纸",sz=13,c=DARK)
ap(tf2,"2️⃣ 用剪刀剪不同形状",sz=13,c=DARK)
ap(tf2,"3️⃣ 在白纸上拼贴",sz=13,c=DARK)
ap(tf2,"4️⃣ 用胶水贴牢",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=GREEN_OK;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"💡 小贴士  Tips",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 形状不用很完美 — 大师也不完美！",sz=13,c=DARK)
ap(tf3,"· 大形状 + 小形状一起拼",sz=13,c=DARK)
ap(tf3,"· 先摆一摆，再贴！",sz=13,c=DARK)
pn(s,n)

# 38 Project 低 Level — Van Gogh sunflower
s=ns();n+=1;bg(s,CREAM);hb(s,"🌻 低 Level: 梵高向日葵  Van Gogh Sunflowers",VANGOGH)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=VANGOGH;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=DARK)
tf=tb(s,0.4,1.45,4.4,1.8,"📄 画纸  Drawing paper",sz=13,c=DARK)
ap(tf,"🖍️ 黄色蜡笔 / 油画棒  Yellow oil pastel",sz=13,c=DARK)
ap(tf,"🟫 棕色蜡笔  Brown crayon",sz=13,c=DARK)
ap(tf,"🎨 水彩 (背景)  Watercolor",sz=13,c=DARK)
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=CORAL;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"🗣️ 句型 Sentence Frames",sz=14,b=True,c=WHITE)
tf_k=tb(s,0.4,3.85,4.4,1.4,"· 这是向日葵。",sz=13,c=DARK)
ap(tf_k,"· 我用了很多黄色。",sz=13,c=DARK)
ap(tf_k,"· 我的向日葵在跳舞。",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=MAGENTA;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 画一个花瓶 (在下面)",sz=13,c=DARK)
ap(tf2,"2️⃣ 画大圆 — 向日葵中心 (棕色)",sz=13,c=DARK)
ap(tf2,"3️⃣ 短粗笔触画花瓣 (黄色)",sz=13,c=DARK)
ap(tf2,"4️⃣ 填颜色 — 像梵高一样有力！",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=GREEN_OK;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"💡 小贴士  Tips",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 笔触要短、要粗、要有力",sz=13,c=DARK)
ap(tf3,"· 多用 2-3 种黄色",sz=13,c=DARK)
ap(tf3,"· 不要怕涂出范围",sz=13,c=DARK)
pn(s,n)

# 39 Project 高 Level — Picasso self-portrait
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 高 Level: 毕加索立体主义自画像  Cubist Self-Portrait",PICASSO)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=PICASSO;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"📄 画纸  Drawing paper",sz=13,c=DARK)
ap(tf,"✏️ 铅笔  Pencil",sz=13,c=DARK)
ap(tf,"🖍️ 彩笔 / 蜡笔  Markers / crayons",sz=13,c=DARK)
ap(tf,"🪞 小镜子  Small mirror",sz=13,c=DARK)
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=GREEN_OK;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"🗣️ 句型 Sentence Frames",sz=14,b=True,c=WHITE)
tf_k=tb(s,0.4,3.85,4.4,1.4,"· 这是我的脸。",sz=13,c=DARK)
ap(tf_k,"· 我用了很多颜色。",sz=13,c=DARK)
ap(tf_k,"· 我的作品很特别。",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=MATISSE;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 照镜子，记住自己的脸",sz=13,c=DARK)
ap(tf2,"2️⃣ 把脸分成 3 块 (横线 + 竖线)",sz=13,c=DARK)
ap(tf2,"3️⃣ 每块用不同角度画 (正面 / 侧面)",sz=13,c=DARK)
ap(tf2,"4️⃣ 加颜色 — 大胆一点！",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=YELLOW;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"💡 小贴士 + 写一句  Tips + 1 line",sz=14,b=True,c=DARK)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 不要怕画得「不像」 — 毕加索就是这样！",sz=13,c=DARK)
ap(tf3,"· 用 4-5 种颜色",sz=13,c=DARK)
ap(tf3,"· 写一句话介绍：「我是____，我喜欢____。」",sz=13,b=True,c=PICASSO)
pn(s,n)

# 40 Gallery walk
s=ns();n+=1;bg(s,CREAM);hb(s,"🖼️ 小小画展  Gallery Walk",YELLOW)
tb(s,0.4,0.9,9,0.5,"把作品贴在墙上，大家一起看！",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Hang your art and walk around like a museum!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("1️⃣","🖼️","贴作品\nHang it up"),
    ("2️⃣","👀","看一看\nLook around"),
    ("3️⃣","🗣️","说一说\nTalk about it"),
    ("4️⃣","⭐","投票\nVote / give stars"),
]
step_colors=[PICASSO,MATISSE,VANGOGH,MAGENTA]
for i,(num,em,d) in enumerate(steps):
    x=0.3+i*2.4;cl=step_colors[i]
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.1,2.0,0.5,num,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.6,2.0,0.8,em,sz=40,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,3.6,2.0,0.4,ls[0],sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# 41 Day 2 Artist Badge
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.4,9,0.8,"🎖️ Day 2 小艺术家徽章  Artist Badge",sz=26,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MAGENTA;sh.line.width=Pt(5)
tf=tb(s,3.6,1.65,2.8,2.7,"DAY 2",sz=18,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"🖼️",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"跟着大师学画画",sz=18,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=13,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
for x,y,cl in [(2.5,2.2,PICASSO),(7.0,2.2,MATISSE),(2.5,3.7,VANGOGH),(7.0,3.7,SKY),(2.0,3.0,GREEN_L),(7.5,3.0,YELLOW)]:
    dot(s,x,y,0.35,cl)
tb(s,1,4.55,8,0.4,"恭喜你！You met 3 masters today! 🎉",sz=16,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"毕加索 · 马蒂斯 · 梵高 · 4 词 · 1 幅作品",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# 42 Tomorrow preview
s=ns();n+=1;bg(s,MAGENTA)
tb(s,0.5,0.9,9,0.8,"🎨 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.1,7,2.7,"Day 3 — 中国水墨画",sz=28,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
ap(tf,"Chinese Ink Painting",sz=16,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🖌️ 毛笔 + 墨 + 水 = ?",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Brush + ink + water = ?",sz=14,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，小艺术家！",sz=15,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/huanli/projects/courseppt/Chinese/小小艺术家little_artist_pbl/day2_masters.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
