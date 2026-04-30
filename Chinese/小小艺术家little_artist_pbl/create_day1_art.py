#!/usr/bin/env python3
"""
小小艺术家 Little Artist Unit — Day 1: 艺术是表达 Art is Expression
Structure modeled on world-trip Day 1, distinct palette (magenta + sunshine + sky).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Playful Art Studio ---
MAGENTA = RGBColor(0xD8,0x1B,0x60)  # primary berry-magenta
YELLOW  = RGBColor(0xFF,0xC1,0x07)  # sunshine
SKY     = RGBColor(0x42,0xA5,0xF5)  # sky blue
CORAL   = RGBColor(0xFF,0x70,0x43)  # coral orange
PURPLE  = RGBColor(0x9C,0x27,0xB0)  # plum
GREEN_L = RGBColor(0x66,0xBB,0x6A)  # leaf green
CREAM   = RGBColor(0xFF,0xF8,0xE7)  # warm cream background
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xFF,0xF3,0xE0)
IMGBG   = RGBColor(0xF0,0xE8,0xF0)   # soft rose-tinted gray
GREEN_OK= RGBColor(0x38,0x8E,0x3C)

# Emotion colors
E_HAPPY = RGBColor(0xFF,0xB3,0x00)  # 开心 bright gold
E_SAD   = RGBColor(0x5C,0x6B,0xC0)  # 难过 indigo
E_ANGRY = RGBColor(0xE5,0x3E,0x3E)  # 生气 red
E_LOVE  = RGBColor(0xEC,0x40,0x7A)  # 喜欢 pink
E_MOOD  = RGBColor(0xAB,0x47,0xBC)  # 心情 purple

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=MAGENTA,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n")
    nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph();p.text=line
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    # playful dots motif — pick 4 colors that contrast with the bg
    palette=[YELLOW,CORAL,SKY,GREEN_L,PURPLE,MAGENTA]
    dot_colors=[c for c in palette if c!=color][:4]
    positions=[(0.8,4.7),(1.6,4.5),(7.8,4.5),(8.6,4.7)]
    for (x,y),cl in zip(positions,dot_colors):
        d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(0.4),Inches(0.4))
        d.fill.solid();d.fill.fore_color.rgb=cl;d.line.fill.background()
    return s

def dot(s,x,y,r,color):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(r),Inches(r))
    d.fill.solid();d.fill.fore_color.rgb=color;d.line.fill.background()

def example_slide(cat_emoji, cat_cn, work_cn, work_en, artist, fact, img_lb, color):
    """Museum-card style: colored header + big image + caption card."""
    s=ns();bg(s,CREAM)
    hb(s,f"{cat_emoji} {cat_cn}  ·  《{work_cn}》",color)
    ib(s,0.5,0.95,9,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,work_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,artist,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def example_slide_generic(cat_emoji, cat_cn, name_cn, name_en, sub_label, fact, img_lb, color):
    """Like example_slide but without 《》 brackets — for instruments, dance types, etc."""
    s=ns();bg(s,CREAM)
    hb(s,f"{cat_emoji} {cat_cn}  ·  {name_cn}",color)
    ib(s,0.5,0.95,9,3.5,img_lb)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.7,4.6,4.5,0.4,name_en,sz=15,b=True,c=color)
    tb(s,0.7,4.95,4.5,0.3,sub_label,sz=11,c=GRAY)
    tb(s,5.3,4.62,4.2,0.75,fact,sz=12,c=DARK)
    return s

def type_overview_slide(emoji, name_cn, name_en, color, hint, subtypes):
    """Category overview — grid of sub-types to introduce before examples."""
    s=ns();bg(s,CREAM);hb(s,f"{emoji} {name_cn}  {name_en}",color)
    tb(s,0.4,0.85,9,0.35,hint,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    cols = 3 if len(subtypes)<=6 else 4
    rows = (len(subtypes)+cols-1)//cols
    card_w = (9.4 - 0.15*(cols-1))/cols
    card_h = 1.7 if rows<=2 else 1.45
    y_start = 1.3 if rows<=2 else 1.25
    y_step = card_h + 0.25 if rows<=2 else card_h + 0.2
    for i,(sem,scn,sen) in enumerate(subtypes):
        col=i%cols;row=i//cols
        x=0.3+col*(card_w+0.15);y=y_start+row*y_step
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(card_w),Inches(card_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
        tb(s,x+0.1,y+0.1,card_w-0.2,0.55,sem,sz=30,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.75,card_w-0.2,0.4,scn,sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+1.15,card_w-0.2,0.3,sen,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    return s

def connect_slide(emoji, art_cn, art_en, color, questions):
    """2-3 connection questions before introducing an art form."""
    s=ns();bg(s,CREAM)
    hb(s,f"{emoji} {art_cn}  {art_en}",color)
    tb(s,0.4,0.9,9.2,0.3,"先聊一聊 — 分享你的经验! / Let's chat — share what YOU know!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    # Big emoji circle on left
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.8),Inches(1.5),Inches(3.0),Inches(3.0))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.8,2.2,3.0,1.6,emoji,sz=110,c=WHITE,a=PP_ALIGN.CENTER)
    # Question cards on right
    nq=len(questions)
    total_h=3.5
    gap=0.15
    card_h=(total_h-gap*(nq-1))/nq
    for i,(q_cn,q_en) in enumerate(questions):
        y=1.5+i*(card_h+gap)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.2),Inches(y),Inches(5.5),Inches(card_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
        nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(4.35),Inches(y+(card_h-0.5)/2),Inches(0.5),Inches(0.5))
        nb.fill.solid();nb.fill.fore_color.rgb=color;nb.line.fill.background()
        tb(s,4.35,y+(card_h-0.5)/2+0.05,0.5,0.4,str(i+1),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,5.0,y+0.13,4.55,0.4,q_cn,sz=14,b=True,c=DARK)
        tb(s,5.0,y+0.55,4.55,0.3,q_en,sz=10,c=GRAY)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.05),Inches(9.4),Inches(0.4))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.4,5.08,9.2,0.32,"💬 没有错答案 — 都说说看! / No wrong answers — everyone share!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    return s

def experience_slide(emoji, art_cn, art_en, color, instruction_cn, instruction_en, students_do, questions, youtube_url):
    """Let's Experience It! — 1-2 min activity with YouTube clip."""
    s=ns();bg(s,CREAM)
    hb(s,f"🎬 一起来试试 {art_cn}!  Let's Experience {art_en}!",color)
    # LEFT: video placeholder + YouTube link
    ib(s,0.3,1.0,4.5,2.85,f"📺 YouTube 视频 / Video clip\n({art_en} activity)")
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.95),Inches(4.5),Inches(0.5))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.4,4.0,4.4,0.4,f"🔗 {youtube_url}",sz=9,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # Big emoji indicator
    tb(s,0.3,4.55,4.5,0.45,emoji+" 1-2 minutes 动起来! Move it!",sz=14,b=True,c=color,a=PP_ALIGN.CENTER)
    # RIGHT — 3 stacked panels
    # 1. Teacher instruction
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(1.0),Inches(4.75),Inches(1.15))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,5.1,1.05,4.5,0.3,"👩‍🏫 老师说明 Teacher Instruction",sz=11,b=True,c=color)
    tb(s,5.1,1.35,4.5,0.35,instruction_cn,sz=13,b=True,c=DARK)
    tb(s,5.1,1.72,4.5,0.35,instruction_en,sz=10,c=GRAY)
    # 2. Students do
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(2.25),Inches(4.75),Inches(1.0))
    sh.fill.solid();sh.fill.fore_color.rgb=YELLOW;sh.line.fill.background()
    tb(s,5.1,2.3,4.5,0.3,"🙋 学生做什么 Students Do",sz=11,b=True,c=DARK)
    tb(s,5.1,2.6,4.5,0.6,students_do,sz=13,b=True,c=DARK)
    # 3. Reflection questions
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(3.35),Inches(4.75),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,5.1,3.4,4.5,0.3,"❓ 试完聊一聊 Reflect",sz=11,b=True,c=color)
    tf=tb(s,5.1,3.7,4.55,0.35,f"❓ {questions[0]}",sz=12,c=DARK)
    for q in questions[1:]:
        ap(tf,"",sz=4)
        ap(tf,f"❓ {q}",sz=12,c=DARK)
    pn=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(5.05),Inches(4.75),Inches(0.4))
    pn.fill.solid();pn.fill.fore_color.rgb=color;pn.line.fill.background()
    tb(s,5.1,5.08,4.5,0.32,"⏱️ 1-2 分钟 / 1-2 minutes",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    return s

def word_card_read(w,py,en,sent,img,color=MAGENTA):
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2)
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=color,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=color)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,img,color=MAGENTA):
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=color)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=color)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    return s

n=0

# ============================================================
# 1 COVER — Artist Palette badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.7,"Little Artist Studio",sz=32,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.45,"小小艺术家",sz=22,c=MAGENTA,a=PP_ALIGN.CENTER)
# big circle badge w/ paint blobs around
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.55),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=MAGENTA;sh.line.color.rgb=YELLOW;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.85),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=YELLOW;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.0,2.8,0.4,"DAY 1",sz=16,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"🎨",sz=56,a=PP_ALIGN.CENTER)
ap(tf,"艺术是表达",sz=20,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"ART IS EXPRESSION",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# paint-splash dots around the badge
for x,y,r,c in [(1.2,1.9,0.6,YELLOW),(1.5,3.8,0.5,SKY),(7.8,1.6,0.55,CORAL),(8.2,3.6,0.5,GREEN_L),(1.0,4.7,0.35,PURPLE),(8.6,4.6,0.35,E_LOVE)]:
    dot(s,x,y,r,c)
tb(s,1,5.05,8,0.4,"🎨 拿起你的画笔，我们出发！Pick up your brush!",sz=14,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","理解艺术是表达 + 认识艺术形式",MAGENTA),
    ("Session 2  下午","2:00-2:45","复习 + 语言目标 (心情词语)",SKY),
    ("Session 3  下午","3:00-4:30","Project：我的心情画 / 我是谁 自画像",CORAL)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=WARM)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"🎨 内容目标  Content:",sz=20,b=True,c=MAGENTA)
tf=tb(s,0.7,1.4,9,1.3,"1. 理解艺术是表达  Art is expression",sz=15,c=DARK)
ap(tf,"2. 认识艺术形式：绘画、音乐、舞蹈……",sz=15,c=DARK)
ap(tf,"3. 发现生活中的艺术，判断什么是艺术",sz=15,c=DARK)
ap(tf,"4. 理解艺术可以表达心情和想法",sz=15,c=DARK)
tb(s,0.5,3.0,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=SKY)
tb(s,0.7,3.5,9,0.4,"👀 我会认：开心  难过  生气  喜欢  心情",sz=15,b=True,c=DARK)
tb(s,0.7,4.0,9,0.4,"✍️ 我会写：开心  喜欢  生气",sz=15,b=True,c=DARK)
tb(s,0.5,4.65,9,0.5,"🖌️ 实践目标：完成 1 幅作品 (心情画 / 自画像)",sz=15,c=CORAL)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","艺术是什么？  What is Art?\n🎨 绘画  🎵 音乐  💃 舞蹈  🎭 戏剧  🎬 电影",MAGENTA,"🖼️")
n+=1

# ============================================================
# 5 WHAT IS ART — big idea slide
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💡 艺术是什么？  What is Art?",MAGENTA)
# Big centered callout
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.2),Inches(1.1),Inches(7.6),Inches(1.6))
sh.fill.solid();sh.fill.fore_color.rgb=MAGENTA;sh.line.fill.background()
tb(s,1.4,1.2,7.2,0.7,"艺术 = 表达",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.4,1.95,7.2,0.7,"Art is Expression",sz=22,c=YELLOW,a=PP_ALIGN.CENTER)
# Three playful bubbles
bubbles=[
    ("🎨","表达看到的","Express what you SEE",YELLOW),
    ("💭","表达想到的","Express what you THINK",SKY),
    ("❤️","表达感受的","Express how you FEEL",CORAL),
]
for i,(em,cn,en,cl) in enumerate(bubbles):
    x=0.5+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.0),Inches(3.0),Inches(2.1))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,3.1,2.8,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.9,2.8,0.45,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.4,2.8,0.35,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 6 ART FORMS — 6-grid
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 艺术形式  Art Forms")
tb(s,0.4,0.9,9,0.4,"艺术有很多种，你认识哪些？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
forms=[
    ("🎨","绘画","Painting",MAGENTA),
    ("🎵","音乐","Music",SKY),
    ("🗿","雕塑","Sculpture",GREEN_L),
    ("🎭","戏剧","Drama",PURPLE),
    ("🎬","电影","Film",YELLOW),
    ("💃","舞蹈","Dance",CORAL),
]
for i,(em,cn,en,cl) in enumerate(forms):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.4+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.1,2.8,0.7,em,sz=40,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.95,2.8,0.45,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.4,2.8,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# PAINTING 绘画 — connect + overview + examples
# ============================================================
s=connect_slide("🎨","绘画","Painting",MAGENTA,[
    ("你喜欢画画吗？","Do you like drawing?"),
    ("你最喜欢画什么？","What do you like to draw?"),
    ("你最喜欢什么颜色？","What's your favorite color?"),
]);n+=1;pn(s,n)

s=type_overview_slide("🎨","绘画","Painting",MAGENTA,
    "绘画有很多种！看看你认识哪些？",
    [("🖼️","油画","Oil Painting"),("🖌️","水墨画","Ink Painting"),
     ("💧","水彩画","Watercolor"),("✏️","素描","Pencil Sketch"),
     ("✂️","拼贴画","Collage"),("🎭","混合媒介","Mixed Media")]);n+=1;pn(s,n)

# 油画 Oil Painting — 蒙娜丽莎
s=example_slide("🖼️","油画","蒙娜丽莎","Mona Lisa",
    "达·芬奇 Leonardo da Vinci · 约 1503 年 · 法国卢浮宫",
    "用油 + 颜料一层一层涂, 颜色慢慢干。\nOil + pigment, layered slowly — most famous smile!",
    "📷 蒙娜丽莎 Mona Lisa",MAGENTA);n+=1;pn(s,n)

# 油画 Oil Painting — 向日葵 (different style)
s=example_slide("🖼️","油画","向日葵","Sunflowers",
    "梵高 Vincent van Gogh · 1888 年 · 荷兰",
    "亮黄色 + 粗笔触 = 像阳光一样温暖！\nBright yellow + bold brush = warm like sunshine!",
    "📷 梵高向日葵 Van Gogh Sunflowers",MAGENTA);n+=1;pn(s,n)

# 水墨画 Ink Painting
s=example_slide("🖌️","水墨画","虾","Shrimp",
    "齐白石 Qi Baishi · 中国 · 近代",
    "毛笔 + 墨 + 水 = 水墨画。只用黑色, 虾像活的!\nBrush + ink + water — alive in just black!",
    "📷 齐白石 虾 Qi Baishi Shrimp",MAGENTA);n+=1;pn(s,n)

# 水彩画 Watercolor
s=example_slide_generic("💧","水彩画 Watercolor","花园一角","A Corner of the Garden",
    "水彩画示例 Watercolor example",
    "用水 + 颜料, 颜色透明像彩虹。轻轻的, 柔柔的。\nTransparent like a rainbow — soft and gentle.",
    "📷 水彩画 Watercolor",MAGENTA);n+=1;pn(s,n)

# 素描 Pencil Sketch
s=example_slide_generic("✏️","素描 Pencil Sketch","铅笔的世界","World of Pencil",
    "只用铅笔 + 橡皮 Just pencil + eraser",
    "深一点、浅一点 — 不用颜色也能画! 线条 + 阴影 = 立体。\nDark or light — no color needed! Lines + shadow = 3D.",
    "📷 铅笔素描 (动物 / 静物 / 人脸)",MAGENTA);n+=1;pn(s,n)

# 拼贴画 Collage — Eric Carle 好饿的毛毛虫
s=example_slide("✂️","拼贴画","好饿的毛毛虫","The Very Hungry Caterpillar",
    "Eric Carle 艾瑞·卡尔 · 1969 年 · 美国绘本作家",
    "用彩纸 + 剪刀 + 胶水 = 故事书里的画！\nColored paper + scissors + glue = picture book art!",
    "📷 Eric Carle 毛毛虫拼贴",MAGENTA);n+=1;pn(s,n)

# 混合媒介 Mixed Media
s=example_slide_generic("🎭","混合媒介 Mixed Media","什么都可以","Anything Goes",
    "现代艺术家 Modern artists",
    "颜料 + 纸 + 布 + 树叶 + 小物品... 全部混在一起!\nPaint + paper + cloth + leaves + objects — all together!",
    "📷 混合媒介作品 (绘画 + 拼贴 + 实物)",MAGENTA);n+=1;pn(s,n)

# ============================================================
# MUSIC 音乐 — connect + overview + examples + experience
# ============================================================
s=connect_slide("🎵","音乐","Music",SKY,[
    ("你喜欢音乐吗？","Do you like music?"),
    ("你最喜欢什么音乐？","What kind of music do you like?"),
    ("你会乐器吗？","Do you play any instrument?"),
]);n+=1;pn(s,n)

s=type_overview_slide("🎵","音乐","Music",SKY,
    "音乐可以用不同的乐器 + 声音！",
    [("🎹","钢琴","Piano"),("🎻","小提琴","Violin"),
     ("🪕","古筝","Guzheng"),("🥁","鼓","Drums"),
     ("🎤","唱歌","Singing"),("🎺","小号","Trumpet")]);n+=1;pn(s,n)

s=example_slide_generic("🎹","钢琴 Piano","《小星星》","Twinkle Twinkle Little Star",
    "莫扎特改编 Mozart variations",
    "钢琴有 88 个键。白键 + 黑键 = 任何歌都能弹！\n88 keys = play any song!",
    "📷 钢琴 Piano",SKY);n+=1;pn(s,n)

s=example_slide_generic("🎻","小提琴 Violin","《梁祝》","Butterfly Lovers",
    "中国著名小提琴协奏曲",
    "小提琴只有 4 根弦，却能唱出故事！\n4 strings tell a whole love story!",
    "📷 小提琴 Violin",SKY);n+=1;pn(s,n)

s=example_slide_generic("🪕","古筝 Guzheng","中国古典音乐","Chinese Classical",
    "中国传统乐器 Traditional Chinese",
    "古筝有 21 根弦，声音像流水一样美。\n21 strings — sounds like flowing water.",
    "📷 古筝 Guzheng",SKY);n+=1;pn(s,n)

s=example_slide_generic("🥁","鼓 Drums","咚咚咚！","Boom Boom Boom!",
    "打击乐 Percussion",
    "鼓声又响又有力，像心跳一样！\nDrums = loud & strong, like a heartbeat!",
    "📷 鼓 Drums",SKY);n+=1;pn(s,n)

# ============================================================
# MUSIC EXPERIENCE — Let's try it!
# ============================================================
s=experience_slide("🎵","音乐","Music",SKY,
    "看视频, 跟着节奏!",
    "Watch the clip and follow the rhythm.",
    "拍手 / 拍腿 / 跟着动!\nClap, tap, or move with the beat.",
    ["简单还是难？  Was it easy or hard?",
     "音乐让你感觉怎么样？  How did the music make you feel?"],
    "youtube.com/watch?v=KUXYMgqw-eg");n+=1;pn(s,n)

# ============================================================
# SCULPTURE 雕塑 — connect + overview + examples
# ============================================================
s=connect_slide("🗿","雕塑","Sculpture",GREEN_L,[
    ("你见过雕像吗？在哪里？","Have you seen sculptures? Where?"),
    ("你玩过黏土 / 橡皮泥吗？","Have you played with clay or play-dough?"),
    ("你想做一个什么雕像？","What sculpture would you make?"),
]);n+=1;pn(s,n)

s=type_overview_slide("🗿","雕塑","Sculpture",GREEN_L,
    "雕塑是立体的艺术，用石头/木头/黏土做！",
    [("🗿","石雕","Stone"),("🪵","木雕","Wood"),
     ("🏺","陶塑","Clay"),("🧊","冰雕","Ice"),
     ("🎨","泥塑","Ceramic"),("🔩","金属","Metal")]);n+=1;pn(s,n)

s=example_slide("🗿","雕塑","摩西像","Moses",
    "米开朗基罗 Michelangelo · 1515 年 · 意大利",
    "大理石雕得像「真人」— 胡子像活的！\nMarble carved so real — even the beard looks alive!",
    "📷 米开朗基罗 摩西像",GREEN_L);n+=1;pn(s,n)

s=example_slide_generic("🏺","雕塑","兵马俑","Terracotta Warriors",
    "中国秦朝 · 2000 多年前",
    "8000 多个泥塑士兵，每一个脸都不一样！\n8,000+ soldiers, each face is different!",
    "📷 兵马俑 Terracotta Warriors",GREEN_L);n+=1;pn(s,n)

# ============================================================
# DRAMA 戏剧 — connect + overview + examples + experience
# ============================================================
s=connect_slide("🎭","戏剧","Drama",PURPLE,[
    ("你知道什么是「戏剧」吗？","Do you know what drama is?"),
    ("你上过台吗？","Have you ever been on stage?"),
    ("能做开心 / 难过 / 生气的脸吗？","Can you make a happy / sad / angry face?"),
]);n+=1;pn(s,n)

s=type_overview_slide("🎭","戏剧","Drama",PURPLE,
    "戏剧 = 演员 + 故事 + 舞台！",
    [("🎼","音乐剧","Musical"),("🎪","话剧","Spoken Drama"),
     ("🎠","木偶戏","Puppet"),("🌑","皮影","Shadow Play"),
     ("🤡","哑剧","Mime"),("🎬","歌剧","Opera")]);n+=1;pn(s,n)

s=example_slide("🎼","音乐剧","狮子王","The Lion King",
    "百老汇音乐剧 · 1997 年首演",
    "唱歌 + 跳舞 + 演戏 = 音乐剧！\nSing + dance + act = musical!",
    "📷 狮子王音乐剧 Lion King",PURPLE);n+=1;pn(s,n)

# ============================================================
# DRAMA EXPERIENCE — Reaction Acting Game
# ============================================================
s=experience_slide("🎭","戏剧","Drama",PURPLE,
    "看演员热身游戏 → 一起玩!",
    "Watch the warm-up game, then play together.",
    "做出表情: 开心 / 难过 / 生气 / 害怕!\nMake faces: happy / sad / angry / scared!",
    ["哪个表情最难做？  Which face was hardest?",
     "演戏好玩吗？  Was acting fun?"],
    "youtube.com/watch?v=zhmcVVTS3mI (Kids Acting Warm-up)");n+=1;pn(s,n)

# ============================================================
# FILM 电影 — connect + overview + examples
# ============================================================
s=connect_slide("🎬","电影","Film",YELLOW,[
    ("你喜欢看电影吗？","Do you like watching movies?"),
    ("你最喜欢的电影是什么？","What's your favorite movie?"),
    ("动画 还是 真人？","Animation or live action?"),
]);n+=1;pn(s,n)

s=type_overview_slide("🎬","电影","Film",YELLOW,
    "电影用镜头讲故事！",
    [("🎞️","动画","Animation"),("🎬","真人","Live Action"),
     ("🦸","超级英雄","Superhero"),("🎠","儿童","Kids"),
     ("🎪","纪录片","Documentary"),("🎭","短片","Short Film")]);n+=1;pn(s,n)

s=example_slide("🎞️","动画电影","哪吒之魔童降世","Ne Zha",
    "中国动画 · 2019 年",
    '"我命由我不由天！" 中国动画很精彩！\nAmazing Chinese animation!',
    "📷 哪吒 Ne Zha",YELLOW);n+=1;pn(s,n)

s=example_slide("🎞️","动画电影","功夫熊猫","Kung Fu Panda",
    "梦工厂 DreamWorks · 2008 年",
    "熊猫也可以当功夫大师！\nEven a panda can be a kung fu master!",
    "📷 功夫熊猫 Kung Fu Panda",YELLOW);n+=1;pn(s,n)

# ============================================================
# DANCE 舞蹈 — connect + overview + examples + experience
# ============================================================
s=connect_slide("💃","舞蹈","Dance",CORAL,[
    ("你喜欢跳舞吗？","Do you like dancing?"),
    ("什么时候你想跳舞？","When do you dance?"),
    ("给我们看一个动作!","Can you show us one dance move?"),
]);n+=1;pn(s,n)

s=type_overview_slide("💃","舞蹈","Dance",CORAL,
    "舞蹈用身体讲故事！",
    [("🩰","芭蕾","Ballet"),("🪭","中国舞","Chinese Dance"),
     ("🕺","街舞","Hip-hop"),("💫","现代舞","Modern"),
     ("🎎","民族舞","Folk"),("🎭","踢踏舞","Tap")]);n+=1;pn(s,n)

s=example_slide("🩰","芭蕾","天鹅湖","Swan Lake",
    "柴可夫斯基 Tchaikovsky · 1877 年",
    "跳舞的人像天鹅一样优雅！\nDancers move like graceful swans!",
    "📷 天鹅湖 Swan Lake ballet",CORAL);n+=1;pn(s,n)

s=example_slide_generic("🪭","中国古典舞","扇子舞","Fan Dance",
    "中国民族舞 Chinese folk",
    "扇子一开一合，像花一样漂亮！\nFans open & close like blooming flowers!",
    "📷 扇子舞 Fan Dance",CORAL);n+=1;pn(s,n)

s=example_slide_generic("🕺","街舞","Hip-hop","Street Dance",
    "现代都市舞蹈 Modern urban",
    "快节奏 + 酷动作 = 自由表达！\nFast beats + cool moves = freedom!",
    "📷 街舞 Hip-hop",CORAL);n+=1;pn(s,n)

# ============================================================
# DANCE EXPERIENCE — Move and Freeze
# ============================================================
s=experience_slide("💃","舞蹈","Dance",CORAL,
    "音乐响 → 动起来! 音乐停 → 定住不动!",
    "Music plays → MOVE!  Music stops → FREEZE!",
    "跟着视频跳 — 用全身!\nDance with the video — use your whole body!",
    ["你定住的动作是什么？  What was your freeze pose?",
     "你最喜欢哪一段？  Which part did you like most?"],
    "youtube.com/watch?v=388Q44ReOWE (Move and Freeze)");n+=1;pn(s,n)

# ============================================================
# 11 ART IN DAILY LIFE — Is this art?
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔍 生活中的艺术  Art in Daily Life",YELLOW)
tb(s,0.4,0.9,9,0.4,"哪些是艺术？Point to the art you see every day!",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
items=[
    ("衣服图案","Clothes patterns",MAGENTA),
    ("漂亮建筑","Buildings",SKY),
    ("摆盘艺术","Food plating",CORAL),
    ("手机壁纸","Phone wallpaper",PURPLE),
    ("绘本插图","Book drawings",GREEN_L),
    ("公园雕塑","Park sculpture",YELLOW),
]
for i,(cn,en,cl) in enumerate(items):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.35+row*2.0
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    # Image placeholder (manual insert)
    ib(s,x+0.15,y+0.1,2.7,0.95,f"📷 {cn}")
    tb(s,x+0.1,y+1.1,2.8,0.4,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.48,2.8,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 12 ART EXPRESSES EMOTIONS
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"❤️ 艺术表达心情  Art Expresses Emotions",E_MOOD)
tb(s,0.4,0.85,9.2,0.32,"同一个心情 — 在不同艺术里，表达方式都不一样!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.28,"Same feeling, different art forms — different ways to show it!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Column headers — 4 art forms
art_forms=[("🎨 绘画","Painting",MAGENTA),("🎵 音乐","Music",SKY),("🎭 戏剧","Drama",PURPLE),("💃 舞蹈","Dance",CORAL)]
col_w=1.95;header_y=1.55;left_w=1.45
for i,(cn,en,cl) in enumerate(art_forms):
    x=left_w+0.1+i*(col_w+0.05)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(header_y),Inches(col_w),Inches(0.55))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,x+0.05,header_y+0.04,col_w-0.1,0.28,cn,sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,header_y+0.31,col_w-0.1,0.22,en,sz=9,c=WARM,a=PP_ALIGN.CENTER)
# 3 emotion rows
emotions=[
    ("😄","开心","HAPPY",E_HAPPY,["亮色\n黄+橙+粉","快节奏\n响亮欢快","笑脸 + 大动作\n跳起来","跳跃 + 转圈"]),
    ("😢","难过","SAD",  E_SAD,  ["蓝色 / 灰色\n冷色调",  "慢, 轻\n钢琴慢调",  "哭脸 + 低头\n慢慢走",  "慢动作 / 蹲下"]),
    ("😡","生气","ANGRY",E_ANGRY,["红色 + 黑色\n粗线条",  "大声 / 鼓\n快重",     "皱眉 + 跺脚\n大声喊",  "用力 + 重重的步"]),
]
row_h=1.05;y0=2.15
for ri,(em,cn,en,e_cl,cells) in enumerate(emotions):
    y=y0+ri*(row_h+0.05)
    # emotion label (left col)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(y),Inches(left_w),Inches(row_h))
    sh.fill.solid();sh.fill.fore_color.rgb=e_cl;sh.line.fill.background()
    txt_c=DARK if e_cl==E_HAPPY else WHITE
    tb(s,0.35,y+0.05,left_w-0.1,0.45,em,sz=28,c=txt_c,a=PP_ALIGN.CENTER)
    tb(s,0.35,y+0.55,left_w-0.1,0.28,cn,sz=14,b=True,c=txt_c,a=PP_ALIGN.CENTER)
    tb(s,0.35,y+0.82,left_w-0.1,0.2,en,sz=9,c=txt_c,a=PP_ALIGN.CENTER)
    # 4 art-form cells
    for ci,cell_text in enumerate(cells):
        x=left_w+0.1+ci*(col_w+0.05)
        col_color=art_forms[ci][2]
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(col_w),Inches(row_h))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=col_color;sh.line.width=Pt(2)
        ls=cell_text.split('\n')
        tf=tb(s,x+0.06,y+0.18,col_w-0.12,0.4,ls[0],sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
        for l in ls[1:]:
            ap(tf,l,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 13 IS THIS ART? — judgment game
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🤔 这是什么艺术？  What Art Form Is This?",YELLOW)
tb(s,0.4,0.85,9.2,0.32,"看图分类 — 它属于哪种艺术？也可以是好几种!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.28,"Sort each picture — what art form? It can be MORE than one!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 6 image cards in 2x3 grid (mix of single-form + multi-form items)
items=[
    ("📷 一幅画 / 蒙娜丽莎","1 种","1 form",MAGENTA),
    ("📷 一首歌 / 钢琴曲","1 种","1 form",SKY),
    ("📷 一座雕像","1 种","1 form",GREEN_L),
    ("📷 动画片 (e.g. 哆啦A梦)","2 种!","2 forms!",YELLOW),
    ("📷 音乐剧《狮子王》","3 种!","3 forms!",PURPLE),
    ("📷 芭蕾舞剧《天鹅湖》","2 种!","2 forms!",CORAL),
]
for i,(img_lb,hint_cn,hint_en,cl) in enumerate(items):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.55+row*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    ib(s,x+0.1,y+0.1,1.95,1.2,img_lb)
    tb(s,x+2.1,y+0.2,0.85,0.4,"?",sz=24,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+2.1,y+0.7,0.85,0.3,hint_cn,sz=11,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+2.1,y+1.0,0.85,0.25,hint_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
# Bottom: 6 art form category badges (the "answer boxes")
tb(s,0.4,4.65,9.2,0.3,"👇 把每张图放进对应的艺术家庭 / Put each into its art family:",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
af=[("🎨 绘画",MAGENTA),("🎵 音乐",SKY),("🗿 雕塑",GREEN_L),("🎭 戏剧",PURPLE),("🎬 电影",YELLOW),("💃 舞蹈",CORAL)]
cell_w=1.5
for i,(label,cl) in enumerate(af):
    x=0.4+i*(cell_w+0.05)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(5.0),Inches(cell_w),Inches(0.4))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,x+0.05,5.04,cell_w-0.1,0.32,label,sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)
notes(s,"老师备课:\n• 让学生指出图片属于哪些艺术家庭, 可以多选!\n• 答案参考:\n   - 一幅画 → 🎨 绘画\n   - 一首歌 → 🎵 音乐\n   - 一座雕像 → 🗿 雕塑\n   - 动画片 → 🎬 电影 + 🎨 绘画\n   - 音乐剧《狮子王》→ 🎵 音乐 + 🎭 戏剧 + 💃 舞蹈\n   - 芭蕾舞剧《天鹅湖》→ 💃 舞蹈 + 🎵 音乐\n• 强调: 艺术可以跨越多种形式, 没有固定边界\n• 鼓励学生举其他例子 — 你最喜欢什么艺术？是不是也是好几种？")

# ============================================================
# 14 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","复习 + 语言目标\n🗣️ 我会认 5 个词  ·  ✍️ 我会写 3 个词",SKY,"📖")
n+=1

# ============================================================
# 15 Review mood match
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔄 快速复习  Quick Review — 连心情",SKY)
tb(s,0.4,0.85,9,0.35,"把表情和心情连起来 (口头)",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
faces=[("😄",E_HAPPY),("😢",E_SAD),("😡",E_ANGRY),("😍",E_LOVE),("🎭",E_MOOD)]
words=["难过","喜欢","心情","开心","生气"]
for i,(em,cl) in enumerate(faces):
    y=1.25+i*0.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(y),Inches(3.4),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.95,y+0.1,3.2,0.5,em,sz=26,c=WHITE,a=PP_ALIGN.CENTER)
for i,w in enumerate(words):
    y=1.25+i*0.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.8),Inches(y),Inches(3.4),Inches(0.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SKY;sh.line.width=Pt(2)
    tb(s,5.95,y+0.12,3.2,0.5,w,sz=22,b=True,c=SKY,a=PP_ALIGN.CENTER)
tb(s,4.25,3.1,1.5,0.5,"?",sz=44,b=True,c=CORAL,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 16-20 我会认 word cards
# ============================================================
read_words=[
    ("开心","kāi xīn","happy","今天我很开心！","📷 开心的脸",E_HAPPY),
    ("难过","nán guò","sad","他难过的时候会画蓝色。","📷 难过的脸",E_SAD),
    ("生气","shēng qì","angry","小朋友生气了，他画了红色。","📷 生气的脸",E_ANGRY),
    ("喜欢","xǐ huān","like","我喜欢画画！","📷 喜欢的脸",E_LOVE),
    ("心情","xīn qíng","mood","画画可以表达我的心情。","📷 各种心情",E_MOOD),
]
for w,py,en,sent,img,cl in read_words:
    s=word_card_read(w,py,en,sent,img,cl);n+=1;pn(s,n)

# ============================================================
# 21 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",SKY)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","老师说词\n学生拍正确的字卡",RGBColor(0xFF,0xF3,0xE0)),
    ("2️⃣","表情配词\nFace Match","看表情\n读对应的词",RGBColor(0xE3,0xF2,0xFD)),
    ("3️⃣","心情接龙\nMood Chain","一个接一个\n说心情词",RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣","我演你猜\nActing","做表情\n大家猜心情",RGBColor(0xFC,0xE4,0xEC)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 22-24 我会写 word cards
# ============================================================
write_words=[
    ("开心","kāi xīn","happy","📷 开心的脸",E_HAPPY),
    ("喜欢","xǐ huān","like","📷 爱心",E_LOVE),
    ("生气","shēng qì","angry","📷 生气的脸",E_ANGRY),
]
for w,py,en,img,cl in write_words:
    s=word_card_write(w,py,en,img,cl);n+=1;pn(s,n)

# ============================================================
# 25 SESSION 3 DIVIDER
# ============================================================
div("Session 3  下午","动手做艺术！  Make Art!\n🎨 我的心情画  ·  👤 我是谁 自画像",CORAL,"🖌️")
n+=1

# ============================================================
# 26 Booklet
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成"艺术是表达"练习册  Day 1 Booklet',CORAL)
ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图 / Booklet pages")
pn(s,n)

# ============================================================
# 27 Projects overview — 3 options (painting + self-portrait + clay sculpture)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 今天做什么？  Pick a Project!",CORAL)
tb(s,0.4,0.9,9,0.4,"3 个选择 — 老师帮你选, 或者你自己选! / Teacher picks, or YOU pick!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
projects=[
    ("低 Level\n(K-2)","🎨 我的心情画","Mood Painting","用颜色画心情\nColors = mood",E_HAPPY,MAGENTA),
    ("高 Level\n(3-5)","👤 我是谁","Self-Portrait","画自己 + 喜欢的东西\nYou + your favs",E_MOOD,PURPLE),
    ("人人可做\n(All Levels)","🟫 粘土雕塑","Clay Sculpture","捏一个小动物 / 心情脸\nMake a small clay piece",GREEN_OK,GREEN_L),
]
card_w=3.0;gap=0.15;start_x=(10-card_w*3-gap*2)/2
for i,(lvl,nm,en,d,ltcl,cl) in enumerate(projects):
    x=start_x+i*(card_w+gap)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(card_w),Inches(3.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(4)
    # level label pill
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(1.55),Inches(card_w-0.4),Inches(0.55))
    sh2.fill.solid();sh2.fill.fore_color.rgb=ltcl;sh2.line.fill.background()
    tf=tb(s,x+0.25,1.6,card_w-0.5,0.3,lvl.split('\n')[0],sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    ap(tf,lvl.split('\n')[1],sz=9,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.25,card_w-0.2,0.5,nm,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.78,card_w-0.2,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.2,3.15,card_w-0.4,1.3,"📷 作品示范")
    ls=d.split('\n')
    tf=tb(s,x+0.1,4.55,card_w-0.2,0.3,ls[0],sz=11,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 28 Project 低 Level — Mood Painting
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 低 Level: 我的心情画  Mood Painting",MAGENTA)
# Left: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=MAGENTA;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"📄 白纸  White paper",sz=13,c=DARK)
ap(tf,"🖍️ 蜡笔 / 彩笔  Crayons / markers",sz=13,c=DARK)
ap(tf,"🎨 水彩 (可选)  Watercolor",sz=13,c=DARK)
ap(tf,"🖼️ 画板  Easel / board",sz=13,c=DARK)
# Middle: color → mood key
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=YELLOW;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"🎨 颜色小贴士  Color Tips",sz=14,b=True,c=DARK)
tf_k=tb(s,0.4,3.85,4.4,1.4,"🟡 黄色 = 开心",sz=13,c=DARK)
ap(tf_k,"🔵 蓝色 = 难过",sz=13,c=DARK)
ap(tf_k,"🔴 红色 = 生气",sz=13,c=DARK)
ap(tf_k,"💗 粉色 = 喜欢",sz=13,c=DARK)
# Right: steps
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SKY;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 想一想：今天我的心情是什么？",sz=13,c=DARK)
ap(tf2,"2️⃣ 选颜色：开心 / 难过 / 生气 / 喜欢",sz=13,c=DARK)
ap(tf2,"3️⃣ 用线条或形状画出心情",sz=13,c=DARK)
ap(tf2,"4️⃣ 可以是抽象的，没有对错！",sz=13,c=DARK)
# Sentence frames
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=GREEN_OK;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"🗣️ 展示句型  Say These",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 这是我的心情画。",sz=13,c=DARK)
ap(tf3,"· 我今天很开心 / 难过。",sz=13,c=DARK)
ap(tf3,"· 我喜欢黄色，因为……",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 29 Project 高 Level — Self-portrait
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👤 高 Level: 我是谁 自画像  Self-Portrait",PURPLE)
# Left: materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=PURPLE;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,1.8,"📄 白纸 / 画板",sz=13,c=DARK)
ap(tf,"✏️ 铅笔 + 彩笔",sz=13,c=DARK)
ap(tf,"🪞 小镜子 (看自己)",sz=13,c=DARK)
ap(tf,"💭 想一想：我喜欢什么？",sz=13,c=DARK)
# Middle: what to include
sh_k=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.35),Inches(4.4),Inches(0.4))
sh_k.fill.solid();sh_k.fill.fore_color.rgb=CORAL;sh_k.line.fill.background()
tb(s,0.4,3.38,4.2,0.35,"💡 可以画什么  Ideas",sz=14,b=True,c=WHITE)
tf_k=tb(s,0.4,3.85,4.4,1.4,"🎨 你喜欢的东西 (宠物、食物、玩具)",sz=13,c=DARK)
ap(tf_k,"🌈 你的心情颜色",sz=13,c=DARK)
ap(tf_k,"⭐ 你擅长什么 (跳舞、画画、运动)",sz=13,c=DARK)
# Right: steps
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=MAGENTA;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,1.8,"1️⃣ 先用铅笔画出自己的脸",sz=13,c=DARK)
ap(tf2,"2️⃣ 加头发、衣服、眼睛",sz=13,c=DARK)
ap(tf2,"3️⃣ 旁边画你喜欢的东西",sz=13,c=DARK)
ap(tf2,"4️⃣ 用颜色表达你的心情",sz=13,c=DARK)
# Sentence frames
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.35),Inches(4.8),Inches(0.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=GREEN_OK;sh3.line.fill.background()
tb(s,5.0,3.38,4.6,0.35,"🗣️ 展示句型  Say These",sz=14,b=True,c=WHITE)
tf3=tb(s,5.0,3.85,4.7,1.4,"· 这是我！",sz=13,c=DARK)
ap(tf3,"· 我喜欢……",sz=13,c=DARK)
ap(tf3,"· 我画了……因为……",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 30 Project — Clay Sculpture (粘土雕塑) — easy options for K-5
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🟫 粘土雕塑  Clay Sculpture (人人可做)",GREEN_L)
# Top-left: Materials
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.6),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=GREEN_L;sh.line.fill.background()
tb(s,0.4,0.98,4.4,0.35,"🧺 材料 Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.6,1.0,"🟫 黏土 (空气干 / 橡皮泥)",sz=12,c=DARK)
ap(tf,"🍢 牙签 (做眼睛、纹路)",sz=12,c=DARK)
ap(tf,"🍽️ 白纸盘 (做底)",sz=12,c=DARK)
ap(tf,"💧 一小杯水 (湿手)",sz=12,c=DARK)
# Top-right: Steps
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(0.95),Inches(4.7),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=CORAL;sh.line.fill.background()
tb(s,5.1,0.98,4.5,0.35,"👉 做法 Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.1,1.45,4.6,1.0,"1️⃣ 揉一揉  Knead the clay soft",sz=12,c=DARK)
ap(tf2,"2️⃣ 搓 / 捏 / 压成基本形状",sz=12,c=DARK)
ap(tf2,"3️⃣ 加细节 (眼睛 / 嘴 / 纹路)",sz=12,c=DARK)
ap(tf2,"4️⃣ 放纸盘上 — 晾干一晚",sz=12,c=DARK)
# Title for ideas
tb(s,0.4,2.55,9.2,0.32,"💡 选一个简单的 — Pick an easy one to make!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
# 6 easy idea cards 2x3 grid
ideas=[
    ("🐢","小乌龟","Turtle","圆 + 4 条腿",GREEN_L),
    ("🐌","蜗牛","Snail","条 + 球 (链 D2 马蒂斯!)",CORAL),
    ("🐱","小猫脸","Cat Face","圆 + 三角耳",MAGENTA),
    ("😊","心情脸","Mood Face","圆 + 表情 (链 D1 心情!)",YELLOW),
    ("🍰","杯子蛋糕","Cupcake","条 + 圆顶",PURPLE),
    ("🍩","甜甜圈","Donut","条 → 圈",SKY),
]
for i,(em,cn,en,hint,cl) in enumerate(ideas):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=2.95+row*1.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(0.95))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.08,y+0.05,0.65,0.85,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.78,y+0.07,2.15,0.32,cn,sz=14,b=True,c=cl)
    tb(s,x+0.78,y+0.4,2.15,0.25,en,sz=10,c=GRAY)
    tb(s,x+0.78,y+0.65,2.15,0.25,hint,sz=9,c=DARK)
# Sentence frames at bottom
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.05),Inches(9.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=GREEN_L;sh.line.width=Pt(2)
tb(s,0.4,5.08,9.2,0.32,"🗣️ 我做的是 ____。它有 ____。/ I made a ____. It has ____.",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 31 Share & Gallery
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🖼️ 小小画展  Gallery Walk",YELLOW)
tb(s,0.4,0.9,9,0.5,"把作品贴在墙上，大家一起看！",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Hang your art and walk around like a museum!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("1️⃣","🖼️","贴作品\nHang it up"),
    ("2️⃣","👀","看同学的\nLook around"),
    ("3️⃣","🗣️","说一说\nTalk about it"),
    ("4️⃣","⭐","贴星星\nGive stars"),
]
for i,(num,em,d) in enumerate(steps):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=[MAGENTA,SKY,CORAL,YELLOW][i];sh.line.width=Pt(3)
    tb(s,x+0.1,2.1,2.0,0.5,num,sz=22,b=True,c=[MAGENTA,SKY,CORAL,YELLOW][i],a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.6,2.0,0.8,em,sz=40,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,3.6,2.0,0.4,ls[0],sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 31 Day 1 Artist Badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.4,9,0.8,"🎖️ Day 1 小艺术家徽章  Artist Badge",sz=26,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
# Big palette-shaped badge: use circle with 6 small color dots
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MAGENTA;sh.line.width=Pt(5)
tf=tb(s,3.6,1.65,2.8,2.7,"DAY 1",sz=18,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"🎨",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"艺术是表达",sz=20,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=13,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
# paint dots around
for x,y,cl in [(2.5,2.2,YELLOW),(7.0,2.2,SKY),(2.5,3.7,CORAL),(7.0,3.7,PURPLE),(2.0,3.0,E_LOVE),(7.5,3.0,GREEN_L)]:
    dot(s,x,y,0.35,cl)
tb(s,1,4.55,8,0.4,"恭喜你完成 Day 1！You are an artist! 🎉",sz=16,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"认识了 6 种艺术形式 · 5 个心情词 · 1 幅作品",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 32 Tomorrow preview
# ============================================================
s=ns();n+=1;bg(s,MAGENTA)
tb(s,0.5,0.9,9,0.8,"🎨 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.5,"Day 2 — 颜色的秘密",sz=28,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
ap(tf,"Secret of Colors",sz=16,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🌈 红 + 黄 = 什么色？",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"What happens when we mix colors?",sz=14,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，小艺术家！",sz=15,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/小小艺术家little_artist_pbl/day1_art_is_expression.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
