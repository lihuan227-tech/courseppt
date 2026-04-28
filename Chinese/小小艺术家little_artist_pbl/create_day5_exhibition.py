#!/usr/bin/env python3
"""
小小艺术家 Little Artist Unit — Day 5: 我手画我心 · 艺术展览日 / Art Exhibition Day
Capstone day — review + group competition + tie-dye/canvas-bag project + 90-min art scroll.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Festival Rainbow + Gold (capstone celebration) ---
MAGENTA   = RGBColor(0xD8,0x1B,0x60)
YELLOW    = RGBColor(0xFF,0xC1,0x07)
SKY       = RGBColor(0x42,0xA5,0xF5)
CORAL     = RGBColor(0xFF,0x70,0x43)
PURPLE    = RGBColor(0x9C,0x27,0xB0)
GREEN_L   = RGBColor(0x66,0xBB,0x6A)
HOT_PINK  = RGBColor(0xFF,0x1B,0x6B)
JET       = RGBColor(0x14,0x14,0x14)
VERMILLION= RGBColor(0xC0,0x39,0x2B)
GOLD      = RGBColor(0xFF,0xD7,0x00)
CREAM     = RGBColor(0xFF,0xF8,0xE7)
WHITE     = RGBColor(0xFF,0xFF,0xFF)
DARK      = RGBColor(0x2C,0x2C,0x2C)
GRAY      = RGBColor(0x88,0x88,0x88)
LGRAY     = RGBColor(0xBB,0xBB,0xBB)
WARM      = RGBColor(0xFF,0xF3,0xE0)
IMGBG     = RGBColor(0xF0,0xE8,0xF0)
GREEN_OK  = RGBColor(0x38,0x8E,0x3C)

RAINBOW = [MAGENTA, CORAL, YELLOW, GREEN_L, SKY, PURPLE]

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=MAGENTA,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def dot(s,x,y,r,color):
    d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(r),Inches(r))
    d.fill.solid();d.fill.fore_color.rgb=color;d.line.fill.background()
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    palette=[YELLOW,CORAL,SKY,GREEN_L,PURPLE,MAGENTA,HOT_PINK,GOLD]
    dot_colors=[c for c in palette if c!=color][:4]
    positions=[(0.8,4.7),(1.6,4.5),(7.8,4.5),(8.6,4.7)]
    for (x,y),cl in zip(positions,dot_colors):
        d=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(0.4),Inches(0.4))
        d.fill.solid();d.fill.fore_color.rgb=cl;d.line.fill.background()
    return s
def rainbow_dots(s,positions=None):
    """Sprinkle rainbow paint dots across a slide."""
    if positions is None:
        positions=[(0.4,0.4),(9.3,0.35),(0.35,5.0),(9.35,5.0),(0.3,2.7),(9.4,2.7)]
    for i,(x,y) in enumerate(positions):
        dot(s,x,y,0.3,RAINBOW[i%len(RAINBOW)])

n=0

# ============================================================
# 1 COVER — Exhibition badge with rainbow + gold
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.25,8,0.7,"Little Artist Studio",sz=28,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,1,0.85,8,0.4,"小小艺术家",sz=18,c=MAGENTA,a=PP_ALIGN.CENTER)
# Big gold-rimmed badge
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.45),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.color.rgb=HOT_PINK;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.7),Inches(3.0),Inches(3.0))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=GOLD;sh2.line.width=Pt(2.5)
tf=tb(s,3.55,1.85,2.9,0.4,"DAY 5",sz=16,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
ap(tf,"🎨🖼️🎉",sz=44,a=PP_ALIGN.CENTER)
ap(tf,"艺术展览日",sz=20,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"我手画我心",sz=14,c=PURPLE,a=PP_ALIGN.CENTER)
ap(tf,"EXHIBITION DAY",sz=10,b=True,c=GRAY,a=PP_ALIGN.CENTER)
# Rainbow paint splashes around the badge
splashes=[(0.9,1.6,0.55,MAGENTA),(1.4,3.6,0.5,YELLOW),(0.7,4.4,0.4,SKY),
          (8.2,1.5,0.55,CORAL),(7.8,3.6,0.5,GREEN_L),(8.6,4.5,0.4,PURPLE),
          (1.9,1.0,0.3,HOT_PINK),(7.6,1.0,0.3,GOLD)]
for x,y,r,c in splashes:
    dot(s,x,y,r,c)
tb(s,1,5.05,8,0.4,"🎨 我们都是小小艺术家！We are all little artists!",sz=14,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE — three sessions
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule",HOT_PINK)
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","回顾本周内容 + 团队比赛",HOT_PINK),
    ("Session 2  下午","2:00-2:45","T-shirt 扎染  /  帆布包 设计",PURPLE),
    ("Session 3  下午","3:00-4:30","🎨 终极任务：我的艺术卷轴 (90 分钟)",GOLD)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    txt_c = DARK if cl==GOLD else WHITE
    sub_c = DARK if cl==GOLD else WARM
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=txt_c)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=sub_c)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=txt_c)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives",HOT_PINK)
tb(s,0.5,0.85,9,0.45,"🎨 内容目标  Content:",sz=18,b=True,c=MAGENTA)
tf=tb(s,0.7,1.3,9,1.4,"1. 复习 5 天学过的艺术：D1 表达 · D2 大师 · D3 水墨 · D4 当代",sz=14,c=DARK)
ap(tf,"2. 在团队中合作创作",sz=14,c=DARK)
ap(tf,"3. 综合使用本周学到的多种艺术方法",sz=14,c=DARK)
ap(tf,"4. 自信展示自己的作品",sz=14,c=DARK)
tb(s,0.5,2.85,9,0.45,"🗣️ 语言目标 — 应用 Apply (本周已学，今天用！):",sz=18,b=True,c=SKY)
tf2=tb(s,0.7,3.3,9,1.2,"· 这是我们的画。",sz=14,b=True,c=DARK)
ap(tf2,"· 我们画了……",sz=14,b=True,c=DARK)
ap(tf2,"· 我们用了……（颜色 / 点 / 线 / 水墨）",sz=14,b=True,c=DARK)
ap(tf2,"· 我喜欢这一部分。",sz=14,b=True,c=DARK)
tb(s,0.5,4.7,9,0.45,"🖌️ 实践目标：1 件 T-shirt OR 帆布包  +  1 幅小组艺术卷轴",sz=14,c=CORAL)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","回顾本周 + 团队比赛！\nReview the Week + Team Competition",HOT_PINK,"🏆")
n+=1

# ============================================================
# 5 WEEK RECAP — 4-day journey
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌈 我们走过的 4 天  Our 4-Day Journey",HOT_PINK)
tb(s,0.4,0.9,9,0.4,"我们一起看看这周学了什么！Look how far we've come!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
days=[
    ("D1","🎨","艺术是表达","Art is Expression","颜色 + 心情",MAGENTA),
    ("D2","👨‍🎨","跟着大师学画画","Learn from Masters","毕加索 · 马蒂斯 · 梵高",SKY),
    ("D3","🖌️","中国水墨画","Chinese Ink Painting","黑白 · 浓淡 · 竹子",JET),
    ("D4","⚫","当代艺术","Contemporary Art","点 · 线 · 重复 · 视错觉",PURPLE),
]
for i,(d,em,cn,en,sub,cl) in enumerate(days):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(2.25),Inches(3.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    # Top stripe
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(2.25),Inches(0.55))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+0.1,1.5,2.05,0.45,d,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.1,2.05,0.9,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.1,2.05,0.45,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.6,2.05,0.35,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.1,2.05,0.9,sub,sz=11,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,5.2,9,0.3,"⭐ 今天我们用上所有学过的方法！",sz=12,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 6 WHAT DID YOU LEARN? — student talk
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💬 你学到了什么？  What Did You Learn?",HOT_PINK)
tb(s,0.4,0.9,9,0.4,"想一想，告诉同学！Think — share with your friend!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
prompts=[
    ("💡","我学到了……","I learned…",MAGENTA),
    ("❤️","我喜欢……","I liked…",CORAL),
    ("🎨","我会画……","I can draw…",GREEN_L),
    ("🌟","我想再试……","I want to try…",PURPLE),
]
for i,(em,cn,en,cl) in enumerate(prompts):
    col=i%2;row=i//2
    x=0.5+col*4.6;y=1.4+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.4),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.15,y+0.2,1.0,1.2,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.2,y+0.2,3.1,0.55,cn,sz=22,b=True,c=cl)
    tb(s,x+1.2,y+0.85,3.1,0.4,en,sz=12,c=GRAY)
    tb(s,x+1.2,y+1.2,3.1,0.4,"👉 你来说！",sz=11,c=cl)
pn(s,n)

# ============================================================
# 7 GROUP COMPETITION INTRO
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 团队大比拼  Team Competition",HOT_PINK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.0),Inches(8.0),Inches(1.4))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,1.2,1.1,7.6,0.7,"分小组 · 4 轮比赛 · 哪组最棒？",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.2,1.85,7.6,0.4,"4 teams · 4 rounds · Which team wins?",sz=15,c=GOLD,a=PP_ALIGN.CENTER)
# Team setup
info=[
    ("👥","3-4 人一组","3-4 students/team",MAGENTA),
    ("🎯","起一个艺术队名","Pick an art team name",SKY),
    ("🏅","答对得 1 分","1 point for correct",GREEN_L),
    ("🎨","总分高 = 冠军组","Highest score wins!",GOLD),
]
for i,(em,cn,en,cl) in enumerate(info):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.7),Inches(2.25),Inches(2.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,2.85,2.05,0.7,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.65,2.05,0.45,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.15,2.05,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    if cl==GOLD:
        tb(s,x+0.1,4.65,2.05,0.3,"🏆",sz=22,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8 ROUND 1 — Master Match
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 Round 1 · 大师配对  Master Match",SKY)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.4),Inches(4.1))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=SKY;sh.line.width=Pt(3)
tb(s,0.6,1.05,4.0,0.45,"📋 怎么玩  How to Play",sz=16,b=True,c=SKY)
tf=tb(s,0.6,1.5,4.0,3.4,"1️⃣ 老师举一张大师的画",sz=14,c=DARK)
ap(tf,"   Teacher shows a master's painting",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"2️⃣ 各组抢答：是哪位大师？",sz=14,c=DARK)
ap(tf,"   Buzz in — which master?",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"3️⃣ 第一个答对 + 1 分",sz=14,c=DARK)
ap(tf,"   First correct = 1 point",sz=11,c=GRAY)
# Right: 3 master cards
masters=[("🎭","毕加索","Picasso"),("✂️","马蒂斯","Matisse"),("🌻","梵高","Van Gogh")]
for i,(em,cn,en) in enumerate(masters):
    y=1.0+i*1.35
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(y),Inches(4.5),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SKY;sh2.line.width=Pt(2)
    tb(s,5.2,y+0.2,0.9,0.8,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,6.2,y+0.15,3.3,0.5,cn,sz=22,b=True,c=SKY)
    tb(s,6.2,y+0.7,3.3,0.4,en,sz=13,c=GRAY)
ib(s,0.4,5.0,4.4,0.0)  # placeholder, very thin
pn(s,n)

# ============================================================
# 9 ROUND 2 — Style Spotter
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔍 Round 2 · 这是哪种艺术？  Style Spotter",GREEN_L)
tb(s,0.4,0.85,9,0.4,"看图 → 抢答这是哪种艺术？  Look — which style is it?",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
styles=[
    ("🖌️","水墨","Chinese Ink",JET),
    ("🎨","油画","Oil Painting",CORAL),
    ("⚫","当代艺术","Contemporary",PURPLE),
    ("✂️","拼贴","Collage",HOT_PINK),
    ("🌀","视错觉","Op Art",SKY),
]
for i,(em,cn,en,cl) in enumerate(styles):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(1.82),Inches(2.3))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.5,1.72,0.7,em,sz=38,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.3,1.72,0.45,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.8,1.72,0.4,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.05,3.2,1.72,0.45,"📷 范例")
ib(s,0.4,3.95,9.2,1.15,"📷 老师举图：这是哪种艺术？")
pn(s,n)

# ============================================================
# 10 ROUND 3 — Word Chain
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔗 Round 3 · 词语接龙  Word Chain",CORAL)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.4),Inches(4.1))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CORAL;sh.line.width=Pt(3)
tb(s,0.6,1.05,4.0,0.45,"📋 怎么玩  How to Play",sz=16,b=True,c=CORAL)
tf=tb(s,0.6,1.5,4.0,3.4,"1️⃣ 老师说 1 个艺术词",sz=14,c=DARK)
ap(tf,"   Teacher says one art word",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"2️⃣ 下一组接 1 个艺术词",sz=14,c=DARK)
ap(tf,"   Next team must say another",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"3️⃣ 5 秒内答出 = 得 1 分",sz=14,c=DARK)
ap(tf,"   Answer in 5s = 1 point",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"4️⃣ 不能重复！No repeats!",sz=14,b=True,c=CORAL)
# Right: word bank cloud
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(0.95),Inches(4.6),Inches(4.1))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=CORAL;sh2.line.width=Pt(2)
tb(s,5.2,1.05,4.4,0.4,"💡 词库参考  Word Bank",sz=15,b=True,c=CORAL,a=PP_ALIGN.CENTER)
words=[("开心",MAGENTA),("画家",SKY),("颜色",CORAL),("形状",GREEN_L),
       ("艺术",PURPLE),("中国",VERMILLION),("竹子",GREEN_OK),("花",HOT_PINK),
       ("鱼",SKY),("熊猫",JET),("点",PURPLE),("线",MAGENTA),
       ("黑白",JET),("对称",SKY),("花纹",CORAL)]
for i,(w,wcl) in enumerate(words):
    col=i%3;row=i//3
    wx=5.25+col*1.45;wy=1.5+row*0.55
    wsh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(wx),Inches(wy),Inches(1.35),Inches(0.45))
    wsh.fill.solid();wsh.fill.fore_color.rgb=WHITE;wsh.line.color.rgb=wcl;wsh.line.width=Pt(1.5)
    tb(s,wx+0.05,wy+0.05,1.25,0.35,w,sz=14,b=True,c=wcl,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 11 ROUND 4 — Charades
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 Round 4 · 表演大师  Master Charades",PURPLE)
tb(s,0.4,0.85,9,0.4,"一组上台表演大师，其他组抢答！One acts, others guess!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
acts=[
    ("🎭","毕加索","Picasso","摆奇怪表情\nMake weird faces",SKY),
    ("✂️","马蒂斯","Matisse","假装剪纸\nPretend to cut paper",HOT_PINK),
    ("🌻","梵高","Van Gogh","旋转画线条\nSwirling brush strokes",YELLOW),
    ("⚫","草间弥生","Yayoi Kusama","点点点点点\nDots, dots, dots!",PURPLE),
]
for i,(em,cn,en,act,cl) in enumerate(acts):
    col=i%2;row=i//2
    x=0.4+col*4.7;y=1.4+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(1.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.15,1.0,1.2,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+1.15,y+0.1,3.3,0.5,cn,sz=20,b=True,c=cl)
    tb(s,x+1.15,y+0.6,3.3,0.35,en,sz=11,c=GRAY)
    ls=act.split('\n')
    tf=tb(s,x+1.15,y+1.0,3.3,0.35,ls[0],sz=12,b=True,c=DARK)
    for l in ls[1:]:ap(tf,l,sz=10,c=GRAY)
pn(s,n)

# ============================================================
# 12 SCOREBOARD
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 计分板  Scoreboard",GOLD)
# Build a 4-round x 5-team table (1 header col + 4 team cols)
ts=s.shapes.add_table(6,5,Inches(0.4),Inches(0.95),Inches(9.2),Inches(4.0));t=ts.table
t.columns[0].width=Inches(2.4)
for c in range(1,5):t.columns[c].width=Inches(1.7)
hd=["回合 Round","🌈 队 1","🎨 队 2","🌟 队 3","🎭 队 4"]
for c,ct in enumerate(hd):
    cl=t.cell(0,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    rn=p.add_run();rn.text=ct;rn.font.name='KaiTi';rn.font.size=Pt(13);rn.font.bold=True;rn.font.color.rgb=WHITE
    cl.fill.solid();cl.fill.fore_color.rgb=HOT_PINK
rounds=[("🎯 R1 · 大师配对",SKY),("🔍 R2 · 这是哪种艺术",GREEN_L),
        ("🔗 R3 · 词语接龙",CORAL),("🎭 R4 · 表演大师",PURPLE),
        ("🏅 总分 Total",GOLD)]
for r,(label,rcl) in enumerate(rounds,start=1):
    cl=t.cell(r,0);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    rn=p.add_run();rn.text=label;rn.font.name='KaiTi';rn.font.size=Pt(12);rn.font.bold=True
    rn.font.color.rgb=DARK if rcl==GOLD else WHITE
    cl.fill.solid();cl.fill.fore_color.rgb=rcl
    for c in range(1,5):
        cell=t.cell(r,c);cell.text="";tf2=cell.text_frame;tf2.word_wrap=True;p2=tf2.paragraphs[0];p2.alignment=PP_ALIGN.CENTER
        rn2=p2.add_run();rn2.text="___";rn2.font.name='KaiTi';rn2.font.size=Pt(18);rn2.font.bold=True;rn2.font.color.rgb=GRAY
        cell.fill.solid();cell.fill.fore_color.rgb=WARM if r%2 else WHITE
tb(s,0.4,5.05,9,0.3,"🏆 老师边玩边填分 · 高分组 = 冠军组！",sz=12,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 13 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","穿在身上的艺术！\nWearable Art — Tie Dye OR Canvas Bag",PURPLE,"👕")
n+=1

# ============================================================
# 14 PROJECT OPTIONS — 2 cards
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 二选一  Pick Your Project",PURPLE)
tb(s,0.4,0.85,9,0.4,"今天每个人选一个！Each student picks one.",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
projs=[
    ("👕","Project 1","T-shirt 设计","Tie Dye 扎染","把白T扎一扎，染一染\n打开 = 漂亮花纹！",HOT_PINK,GOLD),
    ("👜","Project 2","帆布包 设计","Canvas Bag Design","用画笔/印章/颜料\n空白包变艺术品！",PURPLE,SKY),
]
for i,(em,lbl,cn,en,d,cl,acc) in enumerate(projs):
    x=0.4+i*4.7
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(4.5),Inches(3.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(4)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(1.5),Inches(1.5),Inches(0.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=acc;sh2.line.fill.background()
    tb(s,x+0.25,1.55,1.4,0.4,lbl,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,2.1,4.1,1.0,em,sz=68,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,3.25,4.1,0.55,cn,sz=24,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,3.85,4.1,0.4,en,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.2,4.35,4.1,0.4,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 15 PROJECT 1 — Tie Dye intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👕 Project 1: 什么是扎染？  What is Tie Dye?",HOT_PINK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.0),Inches(9),Inches(1.2))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,0.6,1.1,8.8,0.55,"把布扎起来 → 染颜料 → 打开 = 漂亮花纹！",sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.6,1.7,8.8,0.4,"Tie up the cloth → dye → open = beautiful pattern!",sz=14,c=GOLD,a=PP_ALIGN.CENTER)
patterns=[
    ("🌀","螺旋","Spiral",MAGENTA),
    ("🎯","同心圆","Bullseye",SKY),
    ("〰️","条纹","Stripes",GREEN_L),
]
for i,(em,cn,en,cl) in enumerate(patterns):
    x=0.7+i*3.0
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.6),Inches(2.7),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.7,2.5,0.7,em,sz=44,c=cl,a=PP_ALIGN.CENTER)
    ib(s,x+0.25,3.55,2.2,1.1,f"📷 {cn} 范例")
    tb(s,x+0.1,4.7,2.5,0.4,f"{cn} {en}",sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 16 PROJECT 1 — Materials
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧺 Project 1: 材料  Materials",HOT_PINK)
mats=[
    ("👕","白 T 恤","White T-shirt",HOT_PINK),
    ("🎨","染料 (3-4 色)","Dye (3-4 colors)",CORAL),
    ("🔘","橡皮筋","Rubber bands",YELLOW),
    ("🧤","手套","Gloves",GREEN_L),
    ("🛍️","塑料袋","Plastic bag",SKY),
]
for i,(em,cn,en,cl) in enumerate(mats):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(1.82),Inches(3.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.05,1.55,1.72,0.9,em,sz=46,c=cl,a=PP_ALIGN.CENTER)
    ib(s,x+0.1,2.6,1.62,1.4,"📷")
    tb(s,x+0.05,4.05,1.72,0.45,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.55,1.72,0.4,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 17 PROJECT 1 — 4 Steps
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👉 Project 1: 4 步骤  4 Steps",HOT_PINK)
steps=[
    ("1️⃣","🤲","拧 / 扎 / 折","把 T 恤 (3 种基础玩法)",MAGENTA),
    ("2️⃣","🔘","用橡皮筋绑紧","Tie with rubber bands",CORAL),
    ("3️⃣","💧","滴上染料","Drip the dye on",GREEN_L),
    ("4️⃣","⏰","等 30 分钟","洗，晾干 Wash & dry",SKY),
]
for i,(num,em,cn,en,cl) in enumerate(steps):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.25),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.25),Inches(0.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+0.1,1.05,2.05,0.5,num,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.75,2.05,1.0,em,sz=54,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.0,2.05,0.5,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.55,2.05,1.4,en,sz=11,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 18 PROJECT 1 — 3 Folding Techniques
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌀 Project 1: 3 种扎法  3 Folding Styles",HOT_PINK)
folds=[
    ("🌀","螺旋","Spiral","从中间拧成螺旋\nPinch & twist from center",MAGENTA),
    ("🎯","同心圆","Bullseye","从中心一圈一圈绑\nTie rings from center",SKY),
    ("〰️","条纹","Stripes","卷起来 + 横绑\nRoll up + tie across",GREEN_L),
]
for i,(em,cn,en,d,cl) in enumerate(folds):
    x=0.4+i*3.1
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.9),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.1,2.7,0.9,em,sz=58,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.05,2.7,0.5,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.55,2.7,0.35,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.2,2.95,2.5,1.2,f"📷 {cn} 折法")
    ls=d.split('\n')
    tf=tb(s,x+0.1,4.3,2.7,0.3,ls[0],sz=11,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 19 PROJECT 2 — Canvas Bag intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👜 Project 2: 帆布包  Canvas Bag",PURPLE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.0),Inches(9),Inches(1.4))
sh.fill.solid();sh.fill.fore_color.rgb=PURPLE;sh.line.fill.background()
tb(s,0.6,1.1,8.8,0.55,"用画笔 / 印章 / 颜料",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.6,1.65,8.8,0.55,"把空白帆布包 → 变成艺术品！",sz=22,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.6,2.2,8.8,0.4,"Turn a blank canvas bag into a wearable artwork!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ib(s,1.5,2.7,3.2,2.4,"📷 空白帆布包 Blank bag")
sh3=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(4.85),Inches(3.6),Inches(0.5),Inches(0.6))
sh3.fill.solid();sh3.fill.fore_color.rgb=GOLD;sh3.line.fill.background()
ib(s,5.5,2.7,3.2,2.4,"📷 装饰好的帆布包 Decorated")
pn(s,n)

# ============================================================
# 20 PROJECT 2 — Materials
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧺 Project 2: 材料  Materials",PURPLE)
mats2=[
    ("👜","白色帆布包","White canvas bag",PURPLE),
    ("🎨","布料颜料","Fabric paint",CORAL),
    ("🖍️","布料笔","Fabric markers",HOT_PINK),
    ("🟢","印章 + 模板","Stamps + stencils",GREEN_L),
    ("🖌️","画笔","Brushes",SKY),
]
for i,(em,cn,en,cl) in enumerate(mats2):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(1.82),Inches(3.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.05,1.55,1.72,0.9,em,sz=46,c=cl,a=PP_ALIGN.CENTER)
    ib(s,x+0.1,2.6,1.62,1.4,"📷")
    tb(s,x+0.05,4.05,1.72,0.45,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.55,1.72,0.4,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 21 PROJECT 2 — Design Choices
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 Project 2: 设计选择  Design Choices",PURPLE)
tb(s,0.4,0.85,9,0.4,"用本周学过的方法！Use what we learned this week!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
choices=[
    ("D1","😊","心情色","用颜色表达\nColor = mood",MAGENTA),
    ("D2","🌻","大师风格","学一位大师\nLike a master",SKY),
    ("D3","🖌️","水墨主题","竹子 / 熊猫\nInk theme",JET),
    ("D4","⚫","当代点线","重复花纹 / 对称\nDots & symmetry",PURPLE),
]
for i,(d,em,cn,desc,cl) in enumerate(choices):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(2.25),Inches(3.65))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(2.25),Inches(0.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+0.1,1.45,2.05,0.4,d,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.05,2.05,1.0,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.15,2.05,0.45,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    ls=desc.split('\n')
    tf=tb(s,x+0.1,3.7,2.05,0.4,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,5.2,9,0.3,"💡 也可以混合多种方法！Mix them up!",sz=12,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 22 PROJECT 2 — 4 Steps
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👉 Project 2: 4 步骤  4 Steps",PURPLE)
steps2=[
    ("1️⃣","✏️","先用铅笔轻轻画","Sketch with pencil",MAGENTA),
    ("2️⃣","🎯","选你的「艺术方法」","Pick a method (D1-D4)",SKY),
    ("3️⃣","🎨","用颜料 / 笔涂色","Paint / color it in",CORAL),
    ("4️⃣","✍️","加细节 + 签名","Details + sign it!",GOLD),
]
for i,(num,em,cn,en,cl) in enumerate(steps2):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.25),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.25),Inches(0.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    num_c = DARK if cl==GOLD else WHITE
    tb(s,x+0.1,1.05,2.05,0.5,num,sz=22,b=True,c=num_c,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.75,2.05,1.0,em,sz=54,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.0,2.05,0.5,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.6,2.05,1.4,en,sz=12,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 23 SESSION 3 DIVIDER — GOLD background
# ============================================================
s=ns();n+=1;bg(s,GOLD)
tb(s,1,1.1,8,1.0,"🎨 终极任务",sz=48,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
tb(s,1,2.2,8,0.7,"我的艺术卷轴",sz=36,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,1,3.0,8,0.5,"My Art Scroll · 90 minutes",sz=20,c=PURPLE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3),Inches(3.9),Inches(4),Inches(0.7))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,3,4.0,4,0.5,"⏱️ 90 分钟  ·  5 步骤",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# rainbow dots
splash=[(0.7,1.1,0.4,MAGENTA),(8.9,1.1,0.4,SKY),(0.6,4.6,0.4,GREEN_L),
        (8.9,4.6,0.4,PURPLE),(1.4,4.9,0.3,CORAL),(8.2,4.9,0.3,HOT_PINK)]
for x,y,r,c in splash:dot(s,x,y,r,c)
pn(s,n)

# ============================================================
# 24 ACTIVITY OVERVIEW — Goals + Timeline
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 活动目标 + 时间表  Goals & Timeline",GOLD)
# Top: 4 goal callouts
goals=[
    ("1️⃣","自主选主题","Pick your theme",MAGENTA),
    ("2️⃣","尝试不同材料 / 风格","Try styles & materials",SKY),
    ("3️⃣","展示表达","Show & share",CORAL),
    ("4️⃣","投票","Vote!",PURPLE),
]
for i,(num,cn,en,cl) in enumerate(goals):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(2.25),Inches(1.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,2.05,0.45,num,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.55,2.05,0.5,cn,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.05,2.05,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Bottom: timeline strip
tb(s,0.4,2.8,9,0.4,"⏱️ 5 个步骤 / 90 分钟  ·  Timeline",sz=15,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
timeline=[
    ("Step 1","选主题","10'",MAGENTA),
    ("Step 2","设计","10-15'",SKY),
    ("Step 3","创作","40-45'",CORAL),
    ("Step 4","展示+投票","20'",PURPLE),
    ("Step 5","总结","5'",GOLD),
]
for i,(st,cn,t,cl) in enumerate(timeline):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.4),Inches(1.82),Inches(1.7))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    txt_c = DARK if cl==GOLD else WHITE
    tb(s,x+0.05,3.5,1.72,0.4,st,sz=14,b=True,c=txt_c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.95,1.72,0.5,cn,sz=15,b=True,c=txt_c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.55,1.72,0.5,t,sz=18,b=True,c=txt_c,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 25 STEP 1 — Pick a Theme
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"1️⃣ 选主题 (10 分钟)  Pick a Theme",MAGENTA)
tb(s,0.4,0.85,9,0.4,"小组讨论，选 1 个主题！Discuss with your team — pick 1!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
themes=[
    ("🌈","梦想乐园","Dream Park",HOT_PINK),
    ("🐼","动物世界","Animal World",GREEN_L),
    ("🌳","大自然","Nature",GREEN_OK),
    ("🌟","奇妙世界","Magic World",PURPLE),
    ("🏙️","我的城市","My City",SKY),
]
for i,(em,cn,en,cl) in enumerate(themes):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(1.82),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.05,1.6,1.72,1.0,em,sz=58,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.7,1.72,0.5,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.25,1.72,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.7,1.72,0.5,"你来选！",sz=12,b=True,c=cl,a=PP_ALIGN.CENTER)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.55),Inches(9),Inches(0.6))
sh3.fill.solid();sh3.fill.fore_color.rgb=GOLD;sh3.line.fill.background()
tb(s,0.6,4.65,8.8,0.5,"💡 你们可以画开心的、奇怪的、安静的，都可以！",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 26 STEP 2 — Design Intro (3 things to think)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"2️⃣ 设计作品 (10-15 分钟)  Design Your Piece",SKY)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.0),Inches(8),Inches(0.9))
sh.fill.solid();sh.fill.fore_color.rgb=SKY;sh.line.fill.background()
tb(s,1,1.1,8,0.7,"🧠 想 3 件事  Think About 3 Things",sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
items=[
    ("🎨","画什么？","WHAT to draw?","主角 + 场景\nMain subject + scene",MAGENTA),
    ("🖌️","用什么艺术方法？","WHICH method?","D1-D4 学过的\nFrom this week",CORAL),
    ("👥","谁画哪里？","WHO draws where?","分工合作\nDivide the canvas",GREEN_L),
]
for i,(em,cn,en,d,cl) in enumerate(items):
    x=0.4+i*3.1
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.15),Inches(2.9),Inches(2.95))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,2.25,2.7,0.9,em,sz=56,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.25,2.7,0.5,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.8,2.7,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,4.25,2.7,0.4,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 27 STEP 2 — 6 Art Method Options Grid
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 选你的艺术方法  Pick Your Methods",SKY)
tb(s,0.4,0.83,9,0.35,"本周学过的 6 种方法，可以混合！6 methods — mix them up!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
methods=[
    ("😊","D1","颜色表达心情","开心用亮色，难过用深色",MAGENTA),
    ("👀","D2","不同角度画人物","正面+侧面脸放一起 (毕加索)",SKY),
    ("🌻","D2","亮色 + 旋转线条","短笔触表现动感 (梵高)",YELLOW),
    ("✂️","D2","彩纸拼贴形状","剪圆形/方形/叶子 (马蒂斯)",HOT_PINK),
    ("🌿","D3","水墨","黑白灰、浓淡变化",JET),
    ("⚫","D4","点 + 线 + 花纹","重复点点、对称、视错觉",PURPLE),
]
for i,(em,d,cn,desc,cl) in enumerate(methods):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.25+row*1.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.45))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.1,0.9,1.0,em,sz=38,c=cl,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+1.0),Inches(y+0.1),Inches(0.5),Inches(0.35))
    sh2.fill.solid();sh2.fill.fore_color.rgb=cl;sh2.line.fill.background()
    tb(s,x+1.02,y+0.13,0.45,0.3,d,sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+1.55,y+0.1,1.4,0.4,cn,sz=13,b=True,c=cl) if False else tb(s,x+1.0,y+0.5,1.95,0.45,cn,sz=13,b=True,c=cl)
    tb(s,x+1.0,y+0.95,1.95,0.45,desc,sz=10,c=DARK)
# Tip strip
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.55),Inches(9.2),Inches(0.6))
sh3.fill.solid();sh3.fill.fore_color.rgb=GOLD;sh3.line.fill.background()
tb(s,0.5,4.6,9.0,0.25,"💡 可以混合不同风格！Mix different styles!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,4.85,9.0,0.25,"💡 可以一部分用水墨，一部分用彩色！Half ink, half color!",sz=12,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 28 STEP 3 — Creation Time + Teacher Prompts
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"3️⃣ 创作时间 (40-45 分钟)  Creation Time",CORAL)
# Left: teacher tips card
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.0),Inches(4.5),Inches(4.1))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CORAL;sh.line.width=Pt(3)
tb(s,0.45,1.1,4.2,0.5,"💛 给老师的提示  Teacher Tips",sz=16,b=True,c=CORAL)
tf=tb(s,0.45,1.65,4.3,3.4,"✓ 不要求完美",sz=14,c=DARK)
ap(tf,"  Not aiming for perfect",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"✓ 鼓励多颜色",sz=14,c=DARK)
ap(tf,"  Encourage more colors",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"✓ 作品里有故事",sz=14,c=DARK)
ap(tf,"  Each work tells a story",sz=11,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"✓ 每个人都参与",sz=14,c=DARK)
ap(tf,"  Everyone participates",sz=11,c=GRAY)
# Right: prompts card
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(1.0),Inches(4.7),Inches(4.1))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=HOT_PINK;sh2.line.width=Pt(3)
tb(s,5.1,1.1,4.4,0.5,"🗣️ 老师巡回引导问题  Walk-around Prompts",sz=16,b=True,c=HOT_PINK)
tf2=tb(s,5.1,1.7,4.4,0.45,"❓ 这个地方在发生什么？",sz=15,b=True,c=DARK)
ap(tf2,"   What's happening here?",sz=11,c=GRAY)
ap(tf2,"",sz=4)
ap(tf2,"❓ 可以加点颜色 / 图案吗？",sz=15,b=True,c=DARK)
ap(tf2,"   Can we add some color or patterns?",sz=11,c=GRAY)
ap(tf2,"",sz=4)
ap(tf2,"❓ 这里可以用点点吗？",sz=15,b=True,c=DARK)
ap(tf2,"   Can we use dots here?",sz=11,c=GRAY)
ap(tf2,"",sz=4)
ap(tf2,"❓ 你想加水墨还是彩色？",sz=15,b=True,c=DARK)
ap(tf2,"   Ink or color in this part?",sz=11,c=GRAY)
pn(s,n)

# ============================================================
# 29 STEP 4 — Show & Share (sentence frames)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"4️⃣ 展示 + 投票 (20 分钟)  Show & Vote",PURPLE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.0),Inches(4.4),Inches(4.1))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PURPLE;sh.line.width=Pt(3)
tb(s,0.45,1.1,4.1,0.5,"🎤 展示  Each Group · 1 Minute",sz=16,b=True,c=PURPLE)
tf=tb(s,0.45,1.7,4.1,3.3,"· 一组站在前面",sz=13,c=DARK)
ap(tf,"  Stand in front",sz=10,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"· 一起举起卷轴",sz=13,c=DARK)
ap(tf,"  Hold up the scroll together",sz=10,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"· 每个人说 1 句话",sz=13,c=DARK)
ap(tf,"  Each person says 1 sentence",sz=10,c=GRAY)
ap(tf,"",sz=4)
ap(tf,"⏱️ 1 分钟 / 每组",sz=14,b=True,c=PURPLE)
# Right: sentence frames card
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.0),Inches(4.85),Inches(4.1))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=GREEN_L;sh2.line.width=Pt(3)
tb(s,5.0,1.1,4.6,0.5,"💬 句型  Sentence Frames",sz=16,b=True,c=GREEN_OK)
frames=[
    "👉 这是我们的画。",
    "👉 我们画了……",
    "👉 我们用了……（颜色 / 点 / 线 / 水墨）",
    "👉 我喜欢这一部分。",
]
tf2=tb(s,5.0,1.7,4.6,0.5,frames[0],sz=15,b=True,c=DARK)
for f in frames[1:]:
    ap(tf2,"",sz=6)
    ap(tf2,f,sz=15,b=True,c=DARK)
pn(s,n)

# ============================================================
# 30 STEP 4 — Voting + Awards
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🗳️ 投票！  Vote for Your Favorite!",HOT_PINK)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(0.95),Inches(9),Inches(0.7))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,0.6,1.05,8.8,0.5,"每人 1-2 票  ·  你最喜欢哪一幅？",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
awards=[
    ("🌟","最有创意奖","Most Creative",MAGENTA),
    ("🎨","最好看颜色奖","Best Colors",CORAL),
    ("😄","最开心作品奖","Happiest Art",YELLOW),
]
for i,(em,cn,en,cl) in enumerate(awards):
    x=0.5+i*3.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(2.95),Inches(2.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GOLD;sh.line.width=Pt(4)
    sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.85),Inches(1.95),Inches(1.25),Inches(1.25))
    sh2.fill.solid();sh2.fill.fore_color.rgb=GOLD;sh2.line.color.rgb=cl;sh2.line.width=Pt(2)
    tb(s,x+0.85,2.15,1.25,0.85,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.3,2.75,0.5,cn,sz=17,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.85,2.75,0.4,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.65),Inches(9),Inches(0.55))
sh3.fill.solid();sh3.fill.fore_color.rgb=GOLD;sh3.line.fill.background()
tb(s,0.5,4.7,9,0.45,"💡 也可以简单问：你最喜欢哪一幅？",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 31 STEP 5 — Teacher Wrap-up
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"5️⃣ 老师总结 (5 分钟)  Teacher Wrap-up",GOLD)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.0),Inches(9),Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=HOT_PINK;sh.line.fill.background()
tb(s,0.6,1.1,8.8,0.7,"🎉 今天你们都是小小艺术家！",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.6,1.85,8.8,0.5,"You are ALL little artists today!",sz=18,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.6,2.4,8.8,0.5,"🌈 每个人都很棒！Everyone did amazing!",sz=20,b=True,c=WARM,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.2),Inches(9),Inches(1.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=GOLD;sh2.line.width=Pt(4)
tb(s,0.6,3.3,8.8,0.55,"🎨 每一幅画都不一样，",sz=24,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
tb(s,0.6,3.85,8.8,0.55,"这就是艺术！",sz=28,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.6,4.45,8.8,0.45,"Every painting is different — THAT is art!",sz=15,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 32 WEEK BADGE — 5-day completion
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.25,9,0.7,"🎖️ Little Artist Graduate",sz=30,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
tb(s,0.5,0.85,9,0.4,"小小艺术家结业",sz=20,c=MAGENTA,a=PP_ALIGN.CENTER)
# Big medal in center-left
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.6),Inches(1.5),Inches(3.4),Inches(3.4))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.color.rgb=HOT_PINK;sh.line.width=Pt(5)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.85),Inches(1.75),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=GOLD;sh2.line.width=Pt(3)
tf=tb(s,0.9,1.95,2.8,0.4,"🎨",sz=46,a=PP_ALIGN.CENTER)
ap(tf,"小小艺术家",sz=18,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"GRADUATE",sz=12,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
ap(tf,"✓ 5 DAYS",sz=14,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
# Right: 5-day check list
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.3),Inches(1.5),Inches(5.4),Inches(2.4))
sh3.fill.solid();sh3.fill.fore_color.rgb=WHITE;sh3.line.color.rgb=HOT_PINK;sh3.line.width=Pt(2.5)
tb(s,4.45,1.55,5.1,0.4,"🌟 5 天完成清单  5-Day Checklist",sz=14,b=True,c=HOT_PINK)
days_done=[
    ("✓ D1","🎨","艺术是表达",MAGENTA),
    ("✓ D2","👨‍🎨","跟着大师学画画",SKY),
    ("✓ D3","🖌️","中国水墨画",JET),
    ("✓ D4","⚫","当代艺术",PURPLE),
    ("✓ D5","🎉","艺术展览日",HOT_PINK),
]
for i,(c,em,cn,cl) in enumerate(days_done):
    y=1.95+i*0.38
    tb(s,4.5,y,0.6,0.3,c,sz=13,b=True,c=GREEN_OK)
    tb(s,5.1,y,0.4,0.3,em,sz=15)
    tb(s,5.55,y,4.0,0.3,cn,sz=13,b=True,c=cl)
# Stats strip
sh4=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.3),Inches(4.05),Inches(5.4),Inches(1.05))
sh4.fill.solid();sh4.fill.fore_color.rgb=GOLD;sh4.line.fill.background()
tb(s,4.45,4.1,5.1,0.4,"📊 你学了：",sz=13,b=True,c=DARK)
tb(s,4.45,4.45,5.1,0.3,"6 种艺术形式 · 3 位大师 · 中西画法",sz=12,c=DARK)
tb(s,4.45,4.75,5.1,0.3,"5 种当代艺术 · 完成 5 个项目 🎨",sz=12,b=True,c=HOT_PINK)
tb(s,0.5,5.2,9,0.3,"恭喜你！You are an artist now! 🌈",sz=13,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 33 GROUP PHOTO
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📷 大家一起合影！  Group Photo Time!",HOT_PINK)
tb(s,0.4,0.9,9,0.4,"举着卷轴 + T 恤 / 帆布包 — 一起说「艺术！」",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Hold up your scroll + shirt/bag — say \"ART!\"",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,1.0,1.95,8.0,2.9,"📷 全班合影位 — Class group photo")
# rainbow corner dots
splash=[(0.4,4.95,0.3,MAGENTA),(1.0,5.0,0.25,YELLOW),(8.6,5.0,0.25,SKY),
        (9.3,4.95,0.3,GREEN_L),(0.7,4.85,0.2,CORAL),(8.9,4.85,0.2,PURPLE)]
for x,y,r,c in splash:dot(s,x,y,r,c)
tb(s,0.4,5.2,9,0.3,"📸 茄子！Cheese! · 谷雨中文 · GR EDU",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 34 GOODBYE — colorful celebration
# ============================================================
s=ns();n+=1;bg(s,HOT_PINK)
tb(s,1,0.6,8,0.9,"🎨 再见，小小艺术家！",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.55,8,0.5,"Goodbye, Little Artists!",sz=20,c=GOLD,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(2.3),Inches(7),Inches(2.3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GOLD;sh.line.width=Pt(4)
tf=tb(s,1.7,2.45,6.6,0.5,"See you next summer 🎨!",sz=24,b=True,c=HOT_PINK,a=PP_ALIGN.CENTER)
ap(tf,"明年夏天再见！",sz=22,b=True,c=MAGENTA,a=PP_ALIGN.CENTER)
ap(tf,"",sz=6)
ap(tf,"🌈 继续画画 · Keep painting!",sz=15,c=PURPLE,a=PP_ALIGN.CENTER)
ap(tf,"🖌️ 继续表达 · Keep expressing!",sz=15,c=SKY,a=PP_ALIGN.CENTER)
ap(tf,"❤️ 你就是艺术家! · You ARE an artist!",sz=16,b=True,c=CORAL,a=PP_ALIGN.CENTER)
# Rainbow dots scattered
splashes2=[(0.6,1.0,0.4,YELLOW),(8.9,1.0,0.4,SKY),(0.4,3.4,0.35,GREEN_L),
           (9.1,3.4,0.35,GOLD),(0.6,4.8,0.45,CORAL),(8.8,4.8,0.45,PURPLE),
           (1.4,4.95,0.3,YELLOW),(8.0,4.95,0.3,GOLD)]
for x,y,r,c in splashes2:dot(s,x,y,r,c)
pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/小小艺术家little_artist_pbl/day5_exhibition.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
