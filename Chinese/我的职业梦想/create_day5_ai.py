#!/usr/bin/env python3
"""
我的职业梦想 — Day 5: AI 小帮手来了！ (AI Helper Day)
Focus for young learners (K-5):
  • 复习 Days 1-4 — what students learned in the unit
  • Big idea: 谁多了一个 AI 小帮手? (Doctors/Teachers/Artists/Drivers/Chefs)
  • Core message: AI 可以帮人把事情做得更快、更方便
  • Activities: 猜一猜 (whose AI helper?) + 画一画 (my own AI helper) + 说一说
  • Closing: AI 是小帮手, 人类才是主人
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# === Palette — friendly future tech, soft for K-5 ===
PURPLE = RGBColor(0x6A,0x1B,0x9A)   # primary
CYAN   = RGBColor(0x00,0xBC,0xD4)   # cool accent
NEON   = RGBColor(0x00,0xC8,0x67)   # warm green accent
NAVY   = RGBColor(0x1E,0x3A,0x5F)
GOLD   = RGBColor(0xF5,0xA6,0x23)
CORAL  = RGBColor(0xE9,0x5D,0x5D)
SKY    = RGBColor(0x2E,0x86,0xC1)
HELP   = RGBColor(0xC0,0x39,0x2B)   # helpers red (Day 4)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
SOFT   = RGBColor(0xF2,0xE9,0xFB)   # soft purple wash
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
BROWN  = RGBColor(0x6B,0x44,0x23)

# === Helpers ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=PURPLE,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n,c=GRAY): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=c,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame; lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    return s
def sentence_frame_bar(s,t,frame_cn,frame_en,accent=GOLD):
    if t > 4.95: t = 4.95
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid(); sf.fill.fore_color.rgb=WARM; sf.line.color.rgb=accent; sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.7,0.4,"💬 我来说",sz=14,b=True,c=accent)
    tb(s,2.0,t+0.07,7.6,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.0,t+0.32,7.6,0.3,frame_en,sz=10,c=GRAY)
def mission_card(s,l,t,w,h,num,task_cn,task_en,emoji,color):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=color; sh.line.width=Pt(2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l+0.1),Inches(t+0.08),Inches(0.55),Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb=color; badge.line.fill.background()
    tb(s,l+0.1,t+0.18,0.55,0.4,str(num),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,l+0.7,t+0.15,w-0.8,0.4,task_en,sz=10,c=GRAY)
    tb(s,l+0.05,t+0.85,w-0.1,0.7,emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+1.55,w-0.1,0.4,task_cn,sz=15,b=True,c=color,a=PP_ALIGN.CENTER)

n=0

# ============================================================
# 1. COVER  — AI 小帮手来了！
# ============================================================
s=ns(); bg(s,CREAM)
tb(s,1,0.40,8,0.55,"我的职业梦想 · My Dream Career",sz=24,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.40,"Day 5",sz=18,b=True,c=CYAN,a=PP_ALIGN.CENTER)
ring=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.45),Inches(3.5),Inches(3.5))
ring.fill.solid(); ring.fill.fore_color.rgb=PURPLE; ring.line.color.rgb=CYAN; ring.line.width=Pt(6)
inner=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.75),Inches(2.9),Inches(2.9))
inner.fill.solid(); inner.fill.fore_color.rgb=WHITE; inner.line.color.rgb=NEON; inner.line.width=Pt(2)
tf=tb(s,3.55,1.95,2.9,0.4,"DAY 5",sz=14,b=True,c=CYAN,a=PP_ALIGN.CENTER)
ap(tf,"🤖",sz=72,a=PP_ALIGN.CENTER)
ap(tf,"AI 小帮手来了！",sz=20,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
ap(tf,"AI Helper is Here!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.5,5.15,9,0.4,"✨ AI 可以帮人把事情做得更快、更方便 · AI helps people work faster & easier",sz=13,b=True,c=NEON,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 5 开场 (1 分钟):\n• 「今天 — 我们要认识 AI 小帮手!」\n• 强调: AI 不是要 替代 人, 是 帮助 人\n• 5 天旅程的最后一天 — 后面有展示活动!")

# ============================================================
# 2. 5-DAY JOURNEY MAP (Day 5 highlighted as TODAY)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🗺️ 5 天的职业之旅  Our 5-Day Career Journey",NAVY)
tb(s,0.4,0.85,9.2,0.34,"今天是第 5 天 — AI 小帮手来了！",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.28,"Day 5 — Meet the AI helpers for our favorite jobs!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
days_preview=[
    ("Day 1","认识职业世界","Discover Careers","🌍",NAVY,"8 个职业"),
    ("Day 2","小小科学家","Little Scientists","🔬",SKY,"⭐ 爱迪生"),
    ("Day 3","小小企业家","Little Entrepreneurs","💡",GOLD,"⭐ 乔布斯"),
    ("Day 4","社区小帮手","Community Helpers","🏘",HELP,"老师·医生·厨师"),
    ("Day 5","AI 小帮手","AI Helper","🤖",PURPLE,"⭐ 今天!"),
]
for i,(label,cn,en,em,cl,spotlight) in enumerate(days_preview):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.82),Inches(3.45))
    is_today=(i==4)
    sh.fill.solid(); sh.fill.fore_color.rgb=cl if is_today else WHITE
    sh.line.color.rgb=cl; sh.line.width=Pt(3.5 if is_today else 2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.1),Inches(1.65),Inches(0.55),Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb=WHITE if is_today else cl; badge.line.fill.background()
    tb(s,x+0.1,1.74,0.55,0.4,str(i+1),sz=18,b=True,c=cl if is_today else WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.7,1.7,1.1,0.3,label,sz=11,b=True,c=WHITE if is_today else cl)
    tb(s,x+0.05,2.30,1.72,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,1.72,0.4,cn,sz=14,b=True,c=WHITE if is_today else DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,1.72,0.3,en,sz=10,c=CREAM if is_today else GRAY,a=PP_ALIGN.CENTER)
    sep=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+0.25),Inches(3.85),Inches(1.32),Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb=WHITE if is_today else cl; sep.line.fill.background()
    tb(s,x+0.05,4.00,1.72,0.85,spotlight,sz=11,b=True,c=WHITE if is_today else cl,a=PP_ALIGN.CENTER)
tb(s,0.4,5.18,9.2,0.30,"🤖 今天我们 复习 + 认识 AI 小帮手 + 画一画自己的 AI! Recap → Meet AI → Draw your own!",sz=12,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30 秒过渡:\n• 「我们已经走过 4 天 — 今天是最后一天!」\n• 提醒每一天的明星人物")

# ============================================================
# 3. TODAY'S 2-SESSION OVERVIEW
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🧭 今天的两节课  Today's 2 Sessions",PURPLE)
tb(s,0.4,0.85,9.2,0.42,"🔄 复习 → 🤖 AI → 🎨 大项目  Recap → AI → Big Project",sz=20,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Two sessions today — review, then AI + project!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Session 1 card
s1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.70),Inches(4.55),Inches(3.30))
s1.fill.solid(); s1.fill.fore_color.rgb=WHITE; s1.line.color.rgb=NAVY; s1.line.width=Pt(3)
s1h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.70),Inches(4.55),Inches(0.55))
s1h.fill.solid(); s1h.fill.fore_color.rgb=NAVY; s1h.line.fill.background()
tb(s,0.55,1.78,4.30,0.4,"📚 Session 1 · 复习  Review",sz=17,b=True,c=WHITE)
tb(s,0.55,2.40,4.30,0.6,"🎮",sz=44,a=PP_ALIGN.CENTER)
tb(s,0.55,3.10,4.30,0.4,"Day 1-4 知识 大闯关!",sz=18,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,0.55,3.55,4.30,0.28,"Bamboozle Game — Days 1-4 Quiz",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,3.95,4.30,0.32,"✓ 4 天 · 12 个抢答题",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.30,4.30,0.32,"✓ 答对加分 · 分小组比赛",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.65,4.30,0.32,"✓ 谁记得最多 — 谁是冠军!",sz=12,b=True,c=NAVY,a=PP_ALIGN.CENTER)
# Session 2 card
s2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.70),Inches(4.55),Inches(3.30))
s2.fill.solid(); s2.fill.fore_color.rgb=WHITE; s2.line.color.rgb=PURPLE; s2.line.width=Pt(3)
s2h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.70),Inches(4.55),Inches(0.55))
s2h.fill.solid(); s2h.fill.fore_color.rgb=PURPLE; s2h.line.fill.background()
tb(s,5.20,1.78,4.30,0.4,"🤖 Session 2 · AI + 大项目",sz=17,b=True,c=WHITE)
tb(s,5.20,2.40,4.30,0.6,"🎨",sz=44,a=PP_ALIGN.CENTER)
tb(s,5.20,3.10,4.30,0.4,"AI 小帮手 + 我的梦想",sz=18,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,5.20,3.55,4.30,0.28,"AI Helpers + My Dream Project",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,5.20,3.95,4.30,0.32,"✓ AI 是什么？工作变样了！",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,5.20,4.30,4.30,0.32,"✓ 画 我的职业梦想 Poster",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,5.20,4.65,4.30,0.32,"✓ 小组 PPT (AI 配图) + Present!",sz=12,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,5.18,9.2,0.30,"🎯 学完今天 — 你 就 完成 全单元 5 天 课程!",sz=12,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30 秒过渡:\n• 今天 2 节课:\n  - Session 1: Bamboozle 复习游戏 (Day 1-4)\n  - Session 2: AI 小帮手 + 大项目\n• 项目: 我的职业梦想 poster + 小组 PPT (用 AI 配图)\n• 强推荐 bonus: 设计 未来 机器人")

# ============================================================
# 4. SESSION 1 DIVIDER — 复习 + Bamboozle
# ============================================================
s=div("Session 1","🎮 复习 大闯关  Bamboozle Review — Days 1-4",NAVY,"📚"); n+=1; pn(s,n,WHITE)
notes(s,"Bamboozle 设置 (老师):\n• 提前 在 bamboozle.com 输入下面 12 个问题\n• 分 2-4 组 — 抢答\n• 答对 拿 积分 / 转 转盘\n• 每个 Day 3 题, 顺序: Day 1 → 2 → 3 → 4\n• ~15 分钟 一共")

# ============================================================
# 5. DAY 1 REVIEW + 3 BAMBOOZLE QUESTIONS
# ============================================================
day_quizzes=[
    ("Day 1","🌍","职业世界  Career World","兴趣 + 能力 + 帮助 = 职业",NAVY,[
        ("找一份好工作，要看自己的什么？",
         [("A","头发颜色",False),("B","兴趣 + 能力",True),
          ("C","身高","",False),("D","年龄",False)]),
        ("喜欢画画的人, 长大可以当什么？",
         [("A","警察",False),("B","厨师",False),
          ("C","画家 / 设计师",True),("D","司机",False)]),
        ("故事书里, 谁在找自己的梦想职业？",
         [("A","小猫",False),("B","小狗",True),
          ("C","老虎",False),("D","兔子",False)]),
    ]),
    ("Day 2","🔬","小小科学家  Little Scientists","为什么? 怎么办? 不放弃!",SKY,[
        ("谁 发明 了 电灯？",
         [("A","牛顿",False),("B","爱迪生",True),
          ("C","爱因斯坦",False),("D","乔布斯",False)]),
        ("爱迪生 失败 了 多少 次 才 成功？",
         [("A","10 次",False),("B","100 次",False),
          ("C","1000+ 次",True),("D","只 一次",False)]),
        ("科学家 最爱 问 什么 问题？",
         [("A","为什么 / 怎么 知道",True),("B","几点了",False),
          ("C","多少 钱",False),("D","几岁了",False)]),
    ]),
    ("Day 3","💡","小小企业家  Little Entrepreneurs","看见需要 → 做出产品 → 帮助别人",GOLD,[
        ("谁 发明 了 iPhone？",
         [("A","爱迪生",False),("B","乔布斯",True),
          ("C","牛顿",False),("D","马斯克",False)]),
        ("企业家 第一步 要 做 什么？",
         [("A","赚钱",False),("B","看见 需要 / 问题",True),
          ("C","开 公司",False),("D","卖 东西",False)]),
        ("Jobs 说: 「___ 就是 最厉害 的 设计」",
         [("A","贵",False),("B","简单",True),
          ("C","大",False),("D","快",False)]),
    ]),
    ("Day 4","🏘","社区小帮手  Community Helpers","老师·医生·厨师·消防员",HELP,[
        ("老师 在 哪里 工作？",
         [("A","医院",False),("B","学校",True),
          ("C","厨房",False),("D","警察局",False)]),
        ("火灾 来 了 — 找 谁？",
         [("A","老师",False),("B","医生",False),
          ("C","消防员",True),("D","司机",False)]),
        ("谁 帮 我们 看 病？",
         [("A","医生",True),("B","厨师",False),
          ("C","邮递员",False),("D","警察",False)]),
    ]),
]
for day_label,em,topic,phrase,cl,questions in day_quizzes:
    s=ns(); bg(s,CREAM); hb(s,f"🎮 {day_label} 复习 · {topic}",cl)
    # Recap chip
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb=WHITE; chip.line.color.rgb=cl; chip.line.width=Pt(2)
    tb(s,0.55,0.93,1.1,0.4,em,sz=22)
    tb(s,1.65,0.90,5.0,0.30,phrase,sz=14,b=True,c=cl)
    tb(s,1.65,1.18,5.0,0.22,f"What we learned on {day_label}",sz=9,c=GRAY)
    cm=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.0),Inches(0.95),Inches(2.45),Inches(0.35))
    cm.fill.solid(); cm.fill.fore_color.rgb=cl; cm.line.fill.background()
    tb(s,7.0,0.99,2.45,0.30,"🎯 Bamboozle 3 题",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # 3 question cards
    for qi,(qtext,opts) in enumerate(questions):
        x=0.40+qi*3.10
        qcard=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(3.00),Inches(3.50))
        qcard.fill.solid(); qcard.fill.fore_color.rgb=WHITE; qcard.line.color.rgb=cl; qcard.line.width=Pt(2)
        # Q badge
        qbadge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.10),Inches(1.65),Inches(0.85),Inches(0.35))
        qbadge.fill.solid(); qbadge.fill.fore_color.rgb=cl; qbadge.line.fill.background()
        tb(s,x+0.10,1.69,0.85,0.30,f"Q{qi+1}",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        # Question text
        tb(s,x+0.12,2.10,2.80,0.85,qtext,sz=12,b=True,c=DARK)
        # 4 options in 2x2 grid
        for oi,opt in enumerate(opts):
            letter,otext,correct=opt[0],opt[1],opt[2] if len(opt)>2 else False
            ox=x+0.10+(oi%2)*1.42
            oy=3.00+(oi//2)*0.95
            obg=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(ox),Inches(oy),Inches(1.38),Inches(0.80))
            if correct:
                obg.fill.solid(); obg.fill.fore_color.rgb=NEON; obg.line.color.rgb=NEON; obg.line.width=Pt(2)
                tcolor=WHITE
                ltrcolor=WHITE
            else:
                obg.fill.solid(); obg.fill.fore_color.rgb=WARM; obg.line.color.rgb=LGRAY; obg.line.width=Pt(1)
                tcolor=DARK
                ltrcolor=cl
            tb(s,ox+0.05,oy+0.05,0.30,0.30,letter,sz=11,b=True,c=ltrcolor)
            tb(s,ox+0.35,oy+0.05,1.00,0.70,otext,sz=10,b=True,c=tcolor)
            if correct:
                tb(s,ox+1.05,oy+0.40,0.30,0.35,"✓",sz=14,b=True,c=WHITE,a=PP_ALIGN.RIGHT)
    # Footer
    tb(s,0.4,5.20,9.2,0.30,"💡 老师: 把这 3 题 输入 Bamboozle — 答对 +10 分 · Wrong / Right wheel",sz=10,b=True,c=cl,a=PP_ALIGN.CENTER)
    n+=1; pn(s,n)
    notes(s,f"{day_label} 复习 · ~3 分钟:\n• 老师 大声 念 题 + 4 个 选项\n• 学生 抢答 — 或 在 Bamboozle 选\n• 正确答案 已经 用 绿色 ✓ 标出 — 教师 视图 用\n• 答 完 揭晓 → 全班 一起 大声 念 一遍 正确 答案")

# ============================================================
# SESSION 2 DIVIDER — AI + 大项目
# ============================================================
s=div("Session 2","🤖 AI 小帮手 + 🎨 我的职业梦想大项目",PURPLE,"🌟"); n+=1; pn(s,n,WHITE)
notes(s,"过渡:\n• Bamboozle 复习 结束 — 现在 开始 Session 2!\n• Session 2 内容: AI 是什么 + 大项目 (个人 poster + 小组 PPT)")

# ============================================================
# 6. AI 是什么 — 4 cards with EMPTY image placeholders
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"💡 AI 是什么？  What is AI?",PURPLE)
tb(s,0.4,0.85,9.2,0.38,"AI 就像一个 小帮手 — 它会做这 4 件事 ↓",sz=18,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.23,9.2,0.26,"AI is a little helper — it can do these 4 things:",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
defs=[
    ("👀","看","See","看图说话"),
    ("👂","听","Hear","听人说话"),
    ("💬","说","Speak","回答问题"),
    ("🧠","学","Learn","越用越聪明"),
]
for i,(em,cn,en,desc) in enumerate(defs):
    x=0.4+i*2.30
    # Outer card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.2),Inches(3.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=PURPLE; sh.line.width=Pt(2.5)
    # Image placeholder (top, big, clearly marked)
    ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.15),Inches(1.70),Inches(1.90),Inches(2.00))
    ph.fill.solid(); ph.fill.fore_color.rgb=SOFT; ph.line.color.rgb=LGRAY; ph.line.width=Pt(1.25)
    tb(s,x+0.15,2.25,1.90,0.5,"📷",sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,2.78,1.90,0.35,"加图片",sz=12,b=True,c=BROWN,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,3.12,1.90,0.30,"Insert image",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    # Caption (small icon + word below image)
    tb(s,x+0.15,3.80,0.45,0.40,em,sz=22)
    tb(s,x+0.65,3.78,1.45,0.40,cn,sz=18,b=True,c=PURPLE)
    tb(s,x+0.65,4.16,1.45,0.28,en,sz=9,c=GRAY)
    tb(s,x+0.15,4.50,1.95,0.45,desc,sz=11,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3 分钟:\n• AI 不是机器人 — AI 是「住在电脑里 的 小帮手」\n• 老师 准备 真实图片 贴到 4 个 placeholder:\n  - 看: Google 翻译镜头 / 手机相册自动分类\n  - 听: Siri / Alexa\n  - 说: ChatGPT 回答问题 / 智能音箱说话\n  - 学: 抖音/YouTube 推荐 / 游戏 AI 对手\n• 重点: AI 是 帮 人, 不是 替代 人")

# ============================================================
# 7. CAREER MENU — 5 jobs, each gets its own slide next
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🌟 谁多了一个 AI 小帮手？  Who Has an AI Helper?",PURPLE)
tb(s,0.4,0.85,9.2,0.42,"我们认识的 5 个职业 — 每一个 都 多了 一个 AI 小帮手！",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"5 jobs we know — each now has an AI helper!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
careers_menu=[
    ("🩺","医生","Doctor",CORAL),
    ("📚","老师","Teacher",SKY),
    ("🎨","画家","Artist",GOLD),
    ("🚗","司机","Driver",NAVY),
    ("👨‍🍳","厨师","Chef",NEON),
]
for i,(em,cn,en,cl) in enumerate(careers_menu):
    x=0.35+i*1.92
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(1.82),Inches(2.95))
    card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.color.rgb=cl; card.line.width=Pt(2.5)
    tb(s,x+0.05,1.85,1.72,0.85,em,sz=56,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,1.72,0.45,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.20,1.72,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    plus=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.70),Inches(3.55),Inches(0.42),Inches(0.42))
    plus.fill.solid(); plus.fill.fore_color.rgb=PURPLE; plus.line.fill.background()
    tb(s,x+0.70,3.60,0.42,0.32,"+",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.30),Inches(4.10),Inches(1.22),Inches(0.40))
    chip.fill.solid(); chip.fill.fore_color.rgb=PURPLE; chip.line.fill.background()
    tb(s,x+0.30,4.16,1.22,0.30,"🤖 AI 小帮手",sz=10,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.4,4.85,9.2,0.34,"👉 接下来 — 我们 一个一个 看! 看图猜一猜, AI 帮 他 做 什么?",sz=14,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,5.18,9.2,0.26,"Next — let's zoom into each one, with pictures!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"1 分钟 过渡:\n• 「我们 已经 认识 这 5 个 职业 — 每个 都 多了 一个 AI 小帮手!」\n• 「下一页 我们 一个一个 看 — 看 真图片!」\n• 后面 5 页 — 每页 一个 职业 + 3 张 图")

# ============================================================
# 7a-7e. PER-CAREER DETAIL SLIDES (5 slides, each with 3 picture placeholders)
# ============================================================
career_details=[
    ("🩺","医生","Doctor",CORAL,[
        ("👀","看 X 光片","Read X-rays","图：医生看 X 光片"),
        ("💊","算 病人 吃 多少 药","Figure out medicine dose","图：药瓶 / 药丸"),
        ("📋","记 病人 资料","Keep patient notes","图：电脑病历 / 医生写字"),
    ]),
    ("📚","老师","Teacher",SKY,[
        ("✍️","帮 改 作业","Grade homework","图：AI 改作业本"),
        ("📝","出 练习题","Make practice questions","图：学生做题"),
        ("📖","给 学生 讲 故事","Tell stories to kids","图：智能音箱讲故事"),
    ]),
    ("🎨","画家","Artist",GOLD,[
        ("💡","想 画画 点子","Spark drawing ideas","图：AI 生成图"),
        ("🌈","选 漂亮 颜色","Pick pretty colors","图：调色板"),
        ("✏️","画 一张 草稿","Sketch a rough draft","图：AI 草图"),
    ]),
    ("🚗","司机","Driver",NAVY,[
        ("🧭","帮 看 路 / 找 路","Help find the way","图：导航地图"),
        ("🅿️","自动 停车","Park automatically","图：自动停车"),
        ("🚦","提醒 红灯 / 行人","Alert red lights / people","图：智能摄像头"),
    ]),
    ("👨‍🍳","厨师","Chef",NEON,[
        ("🍳","推荐 菜谱","Suggest recipes","图：菜谱手机"),
        ("⏰","提醒 火候","Watch cooking time","图：智能厨房"),
        ("🥗","想 新 菜式","Invent new dishes","图：AI 食物"),
    ]),
]
for em,cn,en,cl,tasks in career_details:
    s=ns(); bg(s,CREAM); hb(s,f"{em} {cn} + AI 小帮手  {en} + AI Helper",cl)
    # Top intro line
    tb(s,0.4,0.85,9.2,0.42,f"🤖 AI 能帮 {cn} 做 这 3 件事 ↓",sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.30,9.2,0.26,f"3 things AI helps the {en.lower()} do:",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # 3 cards horizontally — picture placeholder + task name
    for i,(ticon,tcn,ten,hint) in enumerate(tasks):
        x=0.40+i*3.10
        card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(3.00),Inches(3.20))
        card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.color.rgb=cl; card.line.width=Pt(2.5)
        # number badge (top-left)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.15),Inches(1.78),Inches(0.45),Inches(0.45))
        badge.fill.solid(); badge.fill.fore_color.rgb=cl; badge.line.fill.background()
        tb(s,x+0.15,1.85,0.45,0.32,str(i+1),sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        # picture placeholder (large, centered)
        ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.20),Inches(2.35),Inches(2.60),Inches(1.50))
        ph.fill.solid(); ph.fill.fore_color.rgb=SOFT; ph.line.color.rgb=LGRAY; ph.line.width=Pt(1.25)
        tb(s,x+0.20,2.60,2.60,0.45,"📷",sz=28,a=PP_ALIGN.CENTER)
        tb(s,x+0.20,3.05,2.60,0.30,"加图片",sz=11,b=True,c=BROWN,a=PP_ALIGN.CENTER)
        tb(s,x+0.20,3.37,2.60,0.30,hint,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
        # task icon + name below
        tb(s,x+0.20,3.95,0.45,0.45,ticon,sz=22)
        tb(s,x+0.70,3.95,2.20,0.35,tcn,sz=15,b=True,c=DARK)
        tb(s,x+0.70,4.27,2.20,0.30,ten,sz=9,c=GRAY)
    sentence_frame_bar(s,4.95,f"AI 帮 {cn} ___ — 我觉得 ___ 。",f"AI helps the {en.lower()} ___ — I think ___.",accent=cl)
    n+=1; pn(s,n)
    notes(s,f"2 分钟 / 每页:\n• 大声 念 标题: 「{cn} + AI 小帮手!」\n• 一张图 一张图 看 — 「这是 AI 帮 {cn} 做 什么?」\n• 让 孩子 用 「AI 帮 {cn} ___」 句型 说\n• 关键: AI 是 帮 {cn}, 不是 替代 {cn}!")

# ============================================================
# 8. AI 让工作变了 — How AI Changes Our Work
# ============================================================
s=ns(); bg(s,PURPLE)
tb(s,0.5,0.5,9,0.55,"🌟 有了 AI — 工作 会 变成 怎样?",sz=22,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.5,1.05,9,0.32,"With AI — how does work change?",sz=12,c=WHITE,a=PP_ALIGN.CENTER)
# Big sentence card
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(1.50),Inches(8.4),Inches(1.30))
sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=NEON; sh.line.width=Pt(4)
tb(s,0.9,1.65,8.2,0.55,"AI 可以帮人 把事情 做得 更快、更方便！",sz=22,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.9,2.25,8.2,0.45,"AI helps people do things FASTER and EASIER.",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
# 3 columns: changes
tb(s,0.5,3.00,9,0.32,"🔄 变化在 3 个地方：",sz=15,b=True,c=GOLD,a=PP_ALIGN.CENTER)
changes=[
    ("⏩","更快","Faster","以前 1 小时 — 现在 1 分钟",NEON),
    ("🎯","更准","More Precise","AI 不容易 看错 / 算错",CYAN),
    ("🤝","人 + AI","Human + AI","人 想 主意, AI 帮 做",GOLD),
]
for i,(em,cn,en,desc,cl) in enumerate(changes):
    x=0.5+i*3.05
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.45),Inches(2.90),Inches(1.55))
    chip.fill.solid(); chip.fill.fore_color.rgb=WHITE; chip.line.color.rgb=cl; chip.line.width=Pt(2.5)
    tb(s,x+0.05,3.55,2.80,0.5,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.05,2.80,0.32,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.36,2.80,0.25,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.62,2.80,0.40,desc,sz=10,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,5.18,9,0.30,"✨ 但是 — 决定 怎么 做 · 帮谁 · 做什么 — 还是 人 自己 决定！",sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n,WHITE)
notes(s,"3 分钟:\n• 让 孩子 想 一想: 有了 AI, 大人 工作 会 怎样?\n• 重点 3 点:\n  1. 更快 — 以前 几小时 的 活儿 现在 几分钟\n  2. 更准 — AI 看 / 算 不会累, 不容易 错\n  3. 人 + AI 一起 — 人 想 主意, AI 帮 做\n• 但 决定 还是 人 — AI 是 帮手, 不是 老板")

# ============================================================
# 9. PROJECT INTRO — Final Project: 我的 AI 机器人
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🤖 Final Project: 我的 AI 机器人  My AI Robot",PURPLE)
# Driving Question card
dq=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(0.85),Inches(9.20),Inches(1.20))
dq.fill.solid(); dq.fill.fore_color.rgb=WHITE; dq.line.color.rgb=GOLD; dq.line.width=Pt(3)
dqh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(0.85),Inches(9.20),Inches(0.38))
dqh.fill.solid(); dqh.fill.fore_color.rgb=GOLD; dqh.line.fill.background()
tb(s,0.55,0.89,9.0,0.32,"💭 Driving Question  核心问题",sz=12,b=True,c=WHITE)
tb(s,0.55,1.30,9.0,0.35,"如果你能设计一个 AI 机器人, 它会 帮助 谁? 解决 什么 问题?",sz=15,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.55,1.68,9.0,0.28,"If you could design an AI robot, who would it help and what problem would it solve?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3-step overview
proj_steps=[
    ("1","🎨","Step 1","设计","Design","想 点子 + 画 草图","Brainstorm + sketch",GOLD),
    ("2","🛠️","Step 2","制作","Build","用 你选 的 材料 做出来","Build with chosen materials",CYAN),
    ("3","🎤","Step 3","展览","Expo","机器人 展览会 介绍 给大家","Robot expo presentation",NEON),
]
for i,(num,em,sten,cn,en,desc_cn,desc_en,cl) in enumerate(proj_steps):
    x=0.40+i*3.10
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.25),Inches(3.00),Inches(2.70))
    card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.color.rgb=cl; card.line.width=Pt(3)
    hstrip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.25),Inches(3.00),Inches(0.55))
    hstrip.fill.solid(); hstrip.fill.fore_color.rgb=cl; hstrip.line.fill.background()
    nbadge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.15),Inches(2.35),Inches(0.36),Inches(0.36))
    nbadge.fill.solid(); nbadge.fill.fore_color.rgb=WHITE; nbadge.line.fill.background()
    tb(s,x+0.15,2.38,0.36,0.32,num,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.60,2.35,2.30,0.40,sten,sz=15,b=True,c=WHITE)
    tb(s,x+0.05,2.95,2.90,0.65,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.90,0.40,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.05,2.90,0.26,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.35,2.80,0.32,desc_cn,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.65,2.80,0.25,desc_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,5.05,9.2,0.28,"👥 4-5 人 一组  ·  Teams of 4-5",sz=12,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,5.32,9.2,0.22,"AI = 帮助人类解决问题的聪明工具, 不是魔法!",sz=9,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"2 分钟 项目 介绍:\n• 这 是 5 天 课程 的 收尾 项目!\n• 4-5 人 一组\n• 核心 问题: 你 的 AI 机器人 帮 谁? 解决 什么 问题?\n• 重点: AI 不是 魔法 — 是 帮 人 解决 问题 的 聪明 工具\n• 3 步: 设计 → 制作 → 展览")

# ============================================================
# 10. STEP 1 — DESIGN (K-2 vs G3-5 differentiation)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎨 Step 1 · 设计 你的 AI 机器人  Design Your Robot",GOLD)
tb(s,0.4,0.85,9.2,0.36,"先 想 — 再 画! 你的 机器人 帮 谁? 解决 什么?",sz=15,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.4,1.22,9.2,0.26,"First think — then draw! Who does it help? What does it solve?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Left: K-2 simple version
k2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(1.55),Inches(4.55),Inches(3.55))
k2.fill.solid(); k2.fill.fore_color.rgb=WHITE; k2.line.color.rgb=GOLD; k2.line.width=Pt(3)
k2h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(1.55),Inches(4.55),Inches(0.45))
k2h.fill.solid(); k2h.fill.fore_color.rgb=GOLD; k2h.line.fill.background()
tb(s,0.5,1.60,4.30,0.40,"🐣 K-2  简单版  Simple",sz=14,b=True,c=WHITE)
k2items=[
    ("①","机器人 叫 什么 名字?"),
    ("②","它 长 什么样?  (画 出来)"),
    ("③","它 帮助 谁?"),
    ("④","它 会 做 什么?"),
]
for i,(num,cn) in enumerate(k2items):
    yi=2.15+i*0.42
    tb(s,0.55,yi,0.30,0.32,num,sz=14,b=True,c=GOLD)
    tb(s,0.90,yi,3.95,0.32,cn,sz=11,b=True,c=DARK)
# K-2 sentence stems
ss=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(3.95),Inches(4.15),Inches(1.05))
ss.fill.solid(); ss.fill.fore_color.rgb=WARM; ss.line.color.rgb=GOLD; ss.line.width=Pt(1.5)
tb(s,0.65,4.00,3.95,0.28,"💬 句型 (任选 2-3 句):",sz=10,b=True,c=GOLD)
tb(s,0.65,4.28,3.95,0.24,"• 我的机器人叫 ______。",sz=10,b=True,c=DARK)
tb(s,0.65,4.50,3.95,0.24,"• 它 帮助 ______。 它 会 ______。",sz=10,b=True,c=DARK)
tb(s,0.65,4.72,3.95,0.22,"• 它 让 生活 更 ______。",sz=10,c=DARK)
# Right: G3-5 advanced
g3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.55),Inches(4.60),Inches(3.55))
g3.fill.solid(); g3.fill.fore_color.rgb=WHITE; g3.line.color.rgb=NEON; g3.line.width=Pt(3)
g3h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.55),Inches(4.60),Inches(0.45))
g3h.fill.solid(); g3h.fill.fore_color.rgb=NEON; g3h.line.fill.background()
tb(s,5.20,1.60,4.35,0.40,"🌟 G3-5  AI 发明家挑战",sz=14,b=True,c=WHITE)
g3items=[
    ("①","机器人 名称"),
    ("②","使用 场景 (医院/学校/家/太空...)"),
    ("③","解决 的 问题"),
    ("④","AI 如何 「思考」 / 工作"),
    ("⑤","为什么 比 人 更快 / 更安全?"),
]
for i,(num,cn) in enumerate(g3items):
    yi=2.10+i*0.36
    tb(s,5.20,yi,0.30,0.30,num,sz=13,b=True,c=NEON)
    tb(s,5.55,yi,4.00,0.30,cn,sz=11,b=True,c=DARK)
g3ss=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.20),Inches(3.95),Inches(4.30),Inches(1.05))
g3ss.fill.solid(); g3ss.fill.fore_color.rgb=WARM; g3ss.line.color.rgb=NEON; g3ss.line.width=Pt(1.5)
tb(s,5.30,4.00,4.10,0.28,"💬 句型 (写 完整 句子):",sz=10,b=True,c=NEON)
tb(s,5.30,4.28,4.10,0.24,"• 我的 AI 机器人叫 ______, 它 在 ______ 工作。",sz=9,b=True,c=DARK)
tb(s,5.30,4.50,4.10,0.24,"• 它 用 AI 来 ______, 解决 ______ 的 问题。",sz=9,b=True,c=DARK)
tb(s,5.30,4.72,4.10,0.22,"• 它 比 人 更 ______, 因为 ______。",sz=9,c=DARK)
tb(s,0.4,5.20,9.2,0.30,"⏱ 15 分钟 想 + 草图  ·  📦 材料: 纸 · 彩笔 · 贴纸 · googly eyes",sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"15 分钟 设计:\n• 4-5 人 一组 — 但 每个 孩子 想 自己 的 机器人\n• K-2: 4 个 简单 问题 + 句型\n• G3-5: 5 个 进阶 问题 + 完整 句子\n• 材料 准备: 纸, 彩笔, 贴纸, googly eyes, cardboard, pipe cleaners\n• 老师 走 来 走 去, 鼓励 想 法 不 设限")

# ============================================================
# 11. STEP 2 — CHOOSE CATEGORY + BUILD
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🛠️ Step 2 · 选 主题 + 做 出来  Pick a Theme + Build It",CYAN)
tb(s,0.4,0.85,9.2,0.36,"先 选 一个 主题 — 再 决定 怎么 做 出来！",sz=14,b=True,c=CYAN,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.24,"Pick a theme, then choose how to build it!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 8 categories (2 rows x 4)
cats=[
    ("🩺","医疗 机器人","Medical",CORAL),
    ("📚","学习 机器人","Learning",SKY),
    ("🧹","家务 机器人","Chores",GOLD),
    ("🚒","救援 机器人","Rescue",HELP),
    ("🌎","环保 机器人","Earth",NEON),
    ("🐶","宠物 机器人","Pet",BROWN),
    ("🚀","太空 机器人","Space",NAVY),
    ("👵","陪伴 机器人","Companion",PURPLE),
]
tb(s,0.4,1.55,9.2,0.30,"🎯 选 一个 主题  Pick a Theme",sz=12,b=True,c=CYAN,a=PP_ALIGN.CENTER)
for i,(em,cn,en,cl) in enumerate(cats):
    col=i%4; row=i//4
    x=0.40+col*2.35
    y=1.95+row*1.10
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(1.00))
    card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.color.rgb=cl; card.line.width=Pt(2)
    tb(s,x+0.05,y+0.10,0.65,0.65,em,sz=28)
    tb(s,x+0.70,y+0.12,1.45,0.36,cn,sz=12,b=True,c=cl)
    tb(s,x+0.70,y+0.50,1.45,0.30,en,sz=9,c=GRAY)
# Final product options (bottom strip)
tb(s,0.4,4.18,9.2,0.30,"🎨 选 怎么 做  Pick Your Format",sz=12,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
formats=[
    ("✏️","Poster"),("📦","3D 纸盒"),("🧱","LEGO"),("💻","Canva"),("🦈","Pitch"),
]
for i,(em,name) in enumerate(formats):
    x=0.55+i*1.80
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(4.50),Inches(1.65),Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb=WHITE; chip.line.color.rgb=PURPLE; chip.line.width=Pt(1.5)
    tb(s,x+0.08,4.55,0.45,0.45,em,sz=18)
    tb(s,x+0.55,4.58,1.05,0.40,name,sz=11,b=True,c=PURPLE)
tb(s,0.4,5.20,9.2,0.30,"⏱ 20-25 分钟 制作  ·  📦 材料: cardboard · pipe cleaners · googly eyes",sz=11,b=True,c=CYAN,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"20-25 分钟 制作:\n• 学生 选 一个 主题 + 一个 制作 方式\n• 8 类 主题: 医疗/学习/家务/救援/环保/宠物/太空/陪伴\n• 5 种 制作: Poster / 3D 纸盒 / LEGO / Canva / Shark Tank pitch\n• 老师 帮 G3-5 用 Canva, 帮 K-2 剪 纸盒\n• 提醒: 写 名字 + 主题 — 上台 时 容易 介绍")

# ============================================================
# 12. STEP 3 — AI ROBOT EXPO (presentation + rubric)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎤 Step 3 · AI 机器人 展览会  AI Robot Expo!",NEON)
tb(s,0.4,0.85,9.2,0.36,"像 科学展 一样 — 站 在 你的 机器人 旁边, 介绍 给 大家！",sz=14,b=True,c=NEON,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.24,"Like a science fair — stand by your robot and introduce it!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Left: speech frame
frame=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(1.55),Inches(5.80),Inches(3.55))
frame.fill.solid(); frame.fill.fore_color.rgb=WHITE; frame.line.color.rgb=NEON; frame.line.width=Pt(3)
fh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(1.55),Inches(5.80),Inches(0.45))
fh.fill.solid(); fh.fill.fore_color.rgb=NEON; fh.line.fill.background()
tb(s,0.5,1.60,5.55,0.40,"💬 演讲 模板  Speech Frame (1-2 min)",sz=14,b=True,c=WHITE)
splines=[
    ("1","大家好, 这是 我 的 AI 机器人 ______。","Hi, this is my AI robot ___."),
    ("2","它 在 ______ 工作。","It works in ___."),
    ("3","它 帮助 ______。","It helps ___."),
    ("4","它 可以 ______。","It can ___."),
    ("5","它 用 AI 来 ______。","It uses AI to ___."),
    ("6","它 解决 的 问题 是 ______。","The problem it solves is ___."),
]
for i,(num,cn,en) in enumerate(splines):
    yi=2.10+i*0.48
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.50),Inches(yi+0.05),Inches(0.30),Inches(0.30))
    nb.fill.solid(); nb.fill.fore_color.rgb=NEON; nb.line.fill.background()
    tb(s,0.50,yi+0.07,0.30,0.26,num,sz=10,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.90,yi,5.20,0.28,cn,sz=11,b=True,c=DARK)
    tb(s,0.90,yi+0.26,5.20,0.20,en,sz=8,c=GRAY)
# Right: rubric + tip
rub=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.30),Inches(1.55),Inches(3.35),Inches(3.55))
rub.fill.solid(); rub.fill.fore_color.rgb=WHITE; rub.line.color.rgb=GOLD; rub.line.width=Pt(3)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.30),Inches(1.55),Inches(3.35),Inches(0.45))
rh.fill.solid(); rh.fill.fore_color.rgb=GOLD; rh.line.fill.background()
tb(s,6.45,1.60,3.15,0.40,"⭐ 评分 4 颗星  Rubric",sz=13,b=True,c=WHITE)
criteria=[
    ("⭐","有 创意","Creative idea"),
    ("⭐","用途 清楚","Clear purpose"),
    ("⭐","中文 流利","Chinese speaking"),
    ("⭐","展示 努力","Presentation effort"),
]
for i,(st,cn,en) in enumerate(criteria):
    yi=2.15+i*0.55
    tb(s,6.45,yi,0.45,0.40,st,sz=18)
    tb(s,6.90,yi,2.65,0.30,cn,sz=12,b=True,c=DARK)
    tb(s,6.90,yi+0.28,2.65,0.22,en,sz=8,c=GRAY)
# Bonus expo idea
expo=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.45),Inches(4.45),Inches(3.05),Inches(0.55))
expo.fill.solid(); expo.fill.fore_color.rgb=WARM; expo.line.color.rgb=PURPLE; expo.line.width=Pt(1.5)
tb(s,6.55,4.50,2.85,0.26,"🎪 加 分: 邀请 家长!",sz=11,b=True,c=PURPLE)
tb(s,6.55,4.74,2.85,0.22,"Invite parents to the Expo!",sz=8,c=GRAY)
tb(s,0.4,5.20,9.2,0.30,"⏱ 每人 1-2 分钟 — 站 在 自己 的 机器人 旁",sz=11,b=True,c=NEON,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"15-20 分钟 展览:\n• 把 教室 变成 「机器人 博物馆」\n• 学生 把 作品 摆 在 桌上, 站 旁边\n• 老师 / 同学 / 家长 走 过来 — 学生 介绍 (1-2 分钟)\n• 用 6 句 模板 — K-2 选 2-3 句, G3-5 全部 用\n• 评分: 4 颗星 — 创意/用途/中文/努力\n• 加分: 邀请 家长 来 「Robot Expo」 — 像 科学展!")

# ============================================================
# 13. EXTRA — AI ROBOT EXPO DAY (parents)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎪 加分活动 · AI Robot Expo Day  Parent Open House",PURPLE)
# 强推荐 badge
star=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.40),Inches(0.78),Inches(3.20),Inches(0.36))
star.fill.solid(); star.fill.fore_color.rgb=GOLD; star.line.fill.background()
tb(s,3.40,0.83,3.20,0.30,"⭐ 强推荐 · Highly Recommended ⭐",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.40,"开 一场 「机器人 展览会」 — 邀请 家长 来 参观!",sz=15,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.62,9.2,0.24,"Host a Robot Expo — invite parents to come visit!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Left: setup ideas
setup=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(1.95),Inches(4.55),Inches(3.00))
setup.fill.solid(); setup.fill.fore_color.rgb=WHITE; setup.line.color.rgb=PURPLE; setup.line.width=Pt(3)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(1.95),Inches(4.55),Inches(0.45))
sh2.fill.solid(); sh2.fill.fore_color.rgb=PURPLE; sh2.line.fill.background()
tb(s,0.5,2.00,4.30,0.40,"🛠️ 怎么 布置  How to Set Up",sz=14,b=True,c=WHITE)
setup_items=[
    ("①","教室 变 「博物馆」 — 桌子 排 一圈"),
    ("②","每人 一张 桌子 — 摆 机器人 + 名牌"),
    ("③","学生 站 在 自己 的 作品 旁"),
    ("④","家长 / 同学 走过来 听 介绍"),
    ("⑤","老师 拍照 — 做 班级 相册"),
]
for i,(num,cn) in enumerate(setup_items):
    yi=2.55+i*0.45
    tb(s,0.55,yi,0.30,0.32,num,sz=14,b=True,c=PURPLE)
    tb(s,0.90,yi,3.95,0.32,cn,sz=11,b=True,c=DARK)
# Right: why it works
why=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.95),Inches(4.60),Inches(3.00))
why.fill.solid(); why.fill.fore_color.rgb=WHITE; why.line.color.rgb=NEON; why.line.width=Pt(3)
wh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.95),Inches(4.60),Inches(0.45))
wh.fill.solid(); wh.fill.fore_color.rgb=NEON; wh.line.fill.background()
tb(s,5.20,2.00,4.35,0.40,"💡 为什么 棒  Why It Works",sz=13,b=True,c=WHITE)
wpoints=[
    ("🎤","真实 观众 — 孩子 更 认真"),
    ("👀","家长 看到 5 天 学习 的 成果"),
    ("🌟","每个 孩子 都 是 主角"),
    ("📸","留下 美好 回忆"),
]
for i,(em,cn) in enumerate(wpoints):
    yi=2.55+i*0.55
    tb(s,5.15,yi,0.45,0.45,em,sz=22)
    tb(s,5.65,yi+0.05,3.95,0.40,cn,sz=12,b=True,c=DARK)
tb(s,0.4,5.18,9.2,0.32,"✨ AI = 帮助 人类 解决 问题 的 聪明 工具, 不是 魔法!",sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"机器人 展览会 / Parent Open House:\n• 邀请 家长 来 教室 看 学生 作品\n• 像 科学展 一样 布置 — 每人 一桌\n• 学生 站 在 作品 旁 介绍\n• 给 家长 一张 「观众 反馈卡」 — 写 几个 字 鼓励\n• 老师 拍照 — 做 班级 相册 / 发 朋友圈\n• 这 是 5 天 课程 的 最佳 收尾!")

# ============================================================
# 16. UNIT COMPLETE — 5-day journey done!
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎖️ 全单元完成！  Unit Complete!",PURPLE)
tb(s,0.4,0.85,9.2,0.40,"5 天的旅程 — 你 都 走 完 了！",sz=20,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"5 days, 5 themes — you made it!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
done=[
    ("🌍","Day 1","职业世界",NAVY),
    ("🔬","Day 2","小小科学家",SKY),
    ("💡","Day 3","小小企业家",GOLD),
    ("🏘","Day 4","社区小帮手",HELP),
    ("🤖","Day 5","AI 小帮手",PURPLE),
]
for i,(em,d,cn,cl) in enumerate(done):
    x=0.35+i*1.92
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(1.82),Inches(2.10))
    card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.color.rgb=cl; card.line.width=Pt(2.5)
    tb(s,x+0.05,1.75,1.72,0.7,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.45,1.72,0.30,d,sz=11,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,1.72,0.40,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    chk=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.70),Inches(3.18),Inches(0.42),Inches(0.42))
    chk.fill.solid(); chk.fill.fore_color.rgb=NEON; chk.line.fill.background()
    tb(s,x+0.70,3.20,0.42,0.35,"✓",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
close=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(3.95),Inches(8.8),Inches(1.10))
close.fill.solid(); close.fill.fore_color.rgb=PURPLE; close.line.color.rgb=GOLD; close.line.width=Pt(3)
tb(s,0.7,4.10,8.6,0.5,"🌟 长大后, 你想 成为 谁?",sz=22,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.7,4.55,8.6,0.45,"用 你 的 兴趣 + 能力 + AI 小帮手 — 去 帮助 别人！",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.4,5.15,9.2,0.34,"🎉 Use your interests + skills + an AI helper to help others!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"颁奖 5 分钟:\n• 发 「全单元完成」 徽章\n• 全班 合影 — 拿 着 自己 画 的 AI\n• 一句 期许: 「你 的 兴趣 + 能力 + AI 小帮手 = 一个 改变 世界 的 你！」")

# === Save ===
out = os.path.join(os.path.dirname(__file__), "day5_ai.pptx")
prs.save(out)
print(f"✓ Saved: {out}  ({n} slides)")
