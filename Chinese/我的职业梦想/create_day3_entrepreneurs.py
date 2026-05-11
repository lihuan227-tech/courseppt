#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
我的职业梦想 — Day 3: 小小企业家 (Little Entrepreneurs)
3 sessions × 50 min, Workshop Model (5-phase frame per session):
  🔥 Hook (5) → 📚 Mini-Lesson (10) → 🎯 Active Practice (20-25) → 🌱 Apply (5-10) → 🎤 Share & Close (5)
Featured: 乔布斯 Jobs (iPhone story) + 2 kid entrepreneurs (Lily Born · 9-yr popcorn CEO)
Session 3 — 迷你小生意: 4 projects (发现需要 → 设计产品 → 小小推销员 → 帮助学校发明)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# === Palette ===
BIZ    = RGBColor(0xD4,0x8E,0x1F)   # entrepreneur amber/gold
MONEY  = RGBColor(0x2E,0x7D,0x32)   # money green
IDEA   = RGBColor(0xF5,0xC2,0x42)   # idea spark yellow
NAVY   = RGBColor(0x1E,0x3A,0x5F)
LAB    = RGBColor(0xE5,0x3E,0x3E)   # problem red
GREEN  = RGBColor(0x2E,0x7D,0x32)
PURPLE = RGBColor(0x7B,0x1F,0xA2)
RUST   = RGBColor(0xC4,0x52,0x2A)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
BROWN  = RGBColor(0x6B,0x44,0x23)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
OK     = RGBColor(0x38,0x8E,0x3C)
SKY    = RGBColor(0x1F,0x77,0xB4)

# Phase colors
PH_HOOK   = LAB
PH_MINI   = BIZ
PH_ACTIVE = RUST
PH_APPLY  = GREEN
PH_CLOSE  = IDEA

# Per-entrepreneur colors
JOBS    = RGBColor(0x33,0x33,0x33)  # Apple charcoal
LEGO    = RGBColor(0xD3,0x18,0x18)  # LEGO red
MAYUN   = RGBColor(0xFF,0x6A,0x00)  # Alibaba orange

# Team colors (4 teams across the day)
T_RED    = RGBColor(0xE5,0x3E,0x3E)
T_BLUE   = RGBColor(0x1F,0x77,0xB4)
T_GREEN  = RGBColor(0x2E,0x7D,0x32)
T_YELLOW = RGBColor(0xF5,0xA6,0x23)

# === Helpers ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=BIZ,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,l+0.05,t+h/2-0.18,w-0.10,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame; lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    return s
def sentence_frame_bar(s,t,frame_cn,frame_en,accent=IDEA):
    if t > 4.95: t = 4.95
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid(); sf.fill.fore_color.rgb=WARM; sf.line.color.rgb=accent; sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.7,0.4,"💬 我来说",sz=14,b=True,c=accent)
    tb(s,2.0,t+0.07,7.6,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.0,t+0.32,7.6,0.3,frame_en,sz=10,c=GRAY)

def phase_marker(emoji,phase_cn,phase_en,time_min,color,what_cn,what_en):
    s=ns(); bg(s,color)
    tb(s,1,0.85,8,0.7,emoji,sz=80,a=PP_ALIGN.CENTER)
    tb(s,1,1.85,8,0.6,phase_cn,sz=38,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.50,8,0.4,phase_en,sz=16,c=IDEA,a=PP_ALIGN.CENTER)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.5),Inches(3.10),Inches(3.0),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=IDEA; sh.line.fill.background()
    tb(s,3.5,3.18,3.0,0.4,f"⏱  {time_min} 分钟",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,1,4.00,8,0.5,what_cn,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,4.55,8,0.35,what_en,sz=12,c=IDEA,a=PP_ALIGN.CENTER)
    return s

def score_badge(s):
    teams=[("🔴",T_RED),("🔵",T_BLUE),("🟢",T_GREEN),("🟡",T_YELLOW)]
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.40),Inches(0.78),Inches(2.30),Inches(0.32))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=GRAY; sh.line.width=Pt(0.75)
    tb(s,7.45,0.81,0.45,0.28,"🏆",sz=12,a=PP_ALIGN.CENTER)
    for i,(em,cl) in enumerate(teams):
        tb(s,7.85+i*0.45,0.81,0.40,0.28,f"{em}__",sz=10,b=True,c=cl,a=PP_ALIGN.CENTER)

def group_label(s,t=0.78):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.30),Inches(t),Inches(1.80),Inches(0.32))
    sh.fill.solid(); sh.fill.fore_color.rgb=PH_ACTIVE; sh.line.fill.background()
    tb(s,0.30,t+0.03,1.80,0.28,"👥 分组任务",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)

def tpr_strip(s,t,cue_cn,cue_en):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=PH_HOOK; sh.line.fill.background()
    tb(s,0.5,t+0.05,9.0,0.3,f"🙌 TPR · {cue_cn}",sz=14,b=True,c=WHITE)
    tb(s,0.5,t+0.30,9.0,0.25,cue_en,sz=10,c=IDEA)

def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    """Tian-zi-grid (田字格): square box + dashed-look crosshairs + character + optional pinyin label."""
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mid_x=x+size/2; mid_y=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mid_x-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(mid_y-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

n=0

# ============================================================
# 1. COVER
# ============================================================
s=ns(); bg(s,CREAM)
tb(s,1,0.40,8,0.55,"我的职业梦想 · My Dream Career",sz=22,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.4,"Day 3 · 小小企业家  Little Entrepreneurs",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
sh1=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(2.0),Inches(1.55),Inches(2.8),Inches(2.8))
sh1.fill.solid(); sh1.fill.fore_color.rgb=BIZ; sh1.line.color.rgb=IDEA; sh1.line.width=Pt(5)
tf1=tb(s,2.0,1.85,2.8,0.4,"发现需要",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf1,"💡",sz=70,a=PP_ALIGN.CENTER)
ap(tf1,"看见问题!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(5.2),Inches(1.55),Inches(2.8),Inches(2.8))
sh2.fill.solid(); sh2.fill.fore_color.rgb=MONEY; sh2.line.color.rgb=IDEA; sh2.line.width=Pt(5)
tf2=tb(s,5.2,1.85,2.8,0.4,"做出产品",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf2,"🛍️",sz=70,a=PP_ALIGN.CENTER)
ap(tf2,"帮助别人!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,1,4.65,8,0.45,"💰 企业家 = 看见需要 → 做出产品 → 帮助别人",sz=16,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,1,5.15,8,0.3,"3 sessions × 50 min",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30 秒 hook:\n• 「今天三节课，我们都是小老板 — 看见需要、想办法、帮助别人!」\n• 全班分成 4 队: 🔴 红 / 🔵 蓝 / 🟢 绿 / 🟡 黄\n• 准备道具: 假钱 / 贴纸 / 积木 (做 starting capital, S3 用)")

# ============================================================
# 2. SESSION 1 DIVIDER
# ============================================================
s=div("Session 1  上午 11:00–11:50","🌟 故事课  认识乔布斯 + 小朋友老板  ·  50 min",BIZ,"💡"); n+=1; pn(s,n)

# ============================================================
# === SESSION 1 · LEARNING GOALS (objectives overview) ===
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎯 这节课的学习目标  Session 1 Learning Goals",BIZ)
tb(s,0.4,0.85,9.2,0.30,"上完这节课, 你会……",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.15,9.2,0.22,"By the end of this session, you will be able to…",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
goals=[
    ("1","💡","理解什么是「企业家」 — 发现大家的需要, 创造产品或服务, 帮助别人、解决问题。",NAVY),
    ("2","📱","通过乔布斯 iPhone 的故事, 理解企业家的三个步骤: 看见需要 → 做出产品 → 帮助顾客; 体会「简单就是最厉害的设计」和「失败了再改、不放弃」。",JOBS),
    ("3","🦘","通过 Lily 的袋鼠杯(三脚杯)和 9 岁爆米花 CEO 的故事, 明白小孩子也能发现需要、做出产品、解决大问题 — 当企业家不用等长大!",LAB),
    ("4","🌟","说出企业家的共同特点: 看见需要、想点子、不放弃、帮助顾客、勇敢去试。",MONEY),
    ("5","👀","从学校、家、世界、朋友中, 观察并说出自己看到的「需要」。",GREEN),
]
for i,(num,em,text,c) in enumerate(goals):
    y=1.42+i*0.80
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.74))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.13),Inches(0.48),Inches(0.48))
    nb.fill.solid(); nb.fill.fore_color.rgb=c; nb.line.fill.background()
    tb(s,0.55,y+0.18,0.48,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.15,y+0.13,0.55,0.50,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,1.80,y+0.08,7.75,0.62,text,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"1-2 分钟 — 学习目标预告:\n• 老师快速过一遍 5 个目标 — 让学生知道这节课要学什么\n• 不用每条都细讲 — 只要让学生大概知道方向\n• 关键 message: 「这节课结束时, 你会知道企业家是谁、还能像 Lily 一样想自己的小生意!」\n• 引出: 「先从 4 个生活小烦恼开始 — 你也遇到过吗?」")

# ============================================================
# === SESSION 1 · PHASE 1: HOOK (5 min) ===
# ============================================================
# 4. Hook — OPEN question: what problems have you seen? (with example categories)
s=ns(); bg(s,CREAM); hb(s,"🔍 企业家第一步: 发现身边的小问题  Entrepreneur Step 1: Spot Problems Around You",GREEN)
tb(s,0.4,0.85,9.2,0.36,"你有发现过什么问题吗? 你希望可以改善的? 比如……",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.26,"Have you noticed any problems you wish could be fixed? For example…",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 3 category cards (1 row × 3) — EXAMPLES to spark ideas, not to vote on
hook_cats=[
    ("🏫","学校 School",["水壶老漏水","书包太重","铅笔总丢"],SKY),
    ("🏠","家 Home",["玩具玩腻了","东西总找不到","遥控器丢"],PURPLE),
    ("🌍","世界 World",["塑料垃圾多","下雨没伞","老人没人陪"],LAB),
]
for i,(em,label,examples,c) in enumerate(hook_cats):
    x=0.4+i*3.10; y=1.60
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(2.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.15,2.85,0.90,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.05,2.85,0.40,label,sz=16,b=True,c=c,a=PP_ALIGN.CENTER)
    for j,ex in enumerate(examples):
        tb(s,x+0.20,y+1.50+j*0.32,2.65,0.32,f"·  {ex}",sz=12,c=DARK)
# Bottom — OPEN sharing prompt (no vote)
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.30),Inches(9.2),Inches(1.15))
prompt.fill.solid(); prompt.fill.fore_color.rgb=GREEN; prompt.line.color.rgb=IDEA; prompt.line.width=Pt(2.5)
tb(s,0.55,4.38,9.0,0.30,"🙋 轮到你 — 你看到过什么问题? 举手说说看!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,4.70,9.0,0.24,"Your turn — what problems have YOU seen? Hands up and share!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.55,4.98,9.0,0.28,"💬 「我在 ___ 看到过 ___ 的问题。」",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.24,9.0,0.18,"I've seen ___ problem at ___ .",sz=8,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🔥 HOOK · 5-7 分钟 (开放问题 + 例子启发):\n• 老师快速指 4 个格子, 念 1-2 个例子让学生明白「问题」可以是什么样的\n  - 千万不要问「你有没有遇到过这些?」 — 这些只是启发用的例子\n• 然后开放提问: 「你自己 — 在学校、家里、世界、朋友身边 — 看到过什么问题?」\n• 学生举手分享 — 老师把 3-5 个新想法写在白板上 (留下用!)\n• 不评判好坏 — 每一个想法都收下、都点头\n• 句型支持: 「我在 ___ 看到过 ___ 的问题」\n• 老师小结: 「这些问题就是 — 大家的「需要」! 有人需要不漏的水壶, 有人需要陪伴, 有人需要好朋友……」\n• 引出下一页: 「这些问题, 谁来想办法解决? → 企业家!」")

# 5. 答案: 看到 需要 的 人 = 企业家! (parallel to Day 2's Edison reveal)
s=ns(); bg(s,CREAM); hb(s,"💡 答案: 这些都是企业家在做的事!  Answer: That's What Entrepreneurs Do!",BIZ)
tb(s,0.4,0.85,9.2,0.40,"看见这些「需要」,想办法解决 → 这就是企业家!",sz=18,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Seeing these 'needs' and solving them — that's an entrepreneur!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 4 problem → entrepreneur answer cards (compact row)
answers=[
    ("💧","水壶漏","🍶","防漏水壶","Leak-proof bottle",SKY),
    ("🧸","玩具腻","🧱","新款玩具","Cool new toys",PURPLE),
    ("🧥","丢外套","🏷️","姓名贴",   "Name labels",LAB),
    ("🎒","书包重","🎒","轻便书包","Light backpack",BIZ),
]
for i,(p_em,p_cn,a_em,a_cn,a_en,c) in enumerate(answers):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.20),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,1.78,2.10,0.60,p_em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.42,2.10,0.32,p_cn,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.78,2.10,0.40,"↓",sz=18,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.18,2.10,0.60,a_em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.82,2.10,0.32,a_cn,sz=13,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.18,2.10,0.26,a_en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Big takeaway strip
take=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
take.fill.solid(); take.fill.fore_color.rgb=BIZ; take.line.color.rgb=IDEA; take.line.width=Pt(2.5)
tb(s,0.55,4.74,9.0,0.32,"🌟 企业家 = 看见「需要」 → 做出「产品或服务」 → 帮助别人!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.10,9.0,0.26,"Entrepreneur = sees needs → makes product/service → helps people!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-5 分钟 — 答案揭晓:\n• 接 hook: 「这些问题谁来解决? — 企业家!」\n• 一张一张念 4 对:\n  - 水壶漏 → 防漏水壶 (Hydro Flask、S'well 都是这样起步的!)\n  - 玩具腻 → 设计新款玩具 (你也可以!)\n  - 丢外套 → 姓名贴公司 (真的有!)\n  - 书包重 → 轻便设计 (各种新书包品牌)\n• 关键句: 「企业家 = 看见需要 → 做出产品 → 帮助别人」\n• 全班跟读 3 遍 (慢→快→唱)")

# ============================================================
# === SESSION 1 · PHASE 2: MINI-LESSON (10 min) ===
# ============================================================
# === JOBS INTRO 1: Teacher's iPhone — what can it do? ===
s=ns(); bg(s,CREAM); hb(s,"📱 iPhone 可以做什么?  What Can iPhone Do?",JOBS)
tb(s,0.4,0.85,9.2,0.40,"老师拿出自己的 iPhone — 你知道它能做哪些事?",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Teacher shows their iPhone — what can YOU do with one?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 6 things iPhone can do
things=[
    ("📞","打电话","Call",NAVY),
    ("📷","拍照、拍视频","Photo / Video",JOBS),
    ("🎵","听音乐","Music",PURPLE),
    ("🎬","看视频","Watch videos",LAB),
    ("🎮","玩游戏","Games",MONEY),
    ("🗺️","查地图","Maps / GPS",SKY),
]
for i,(em,cn,en,c) in enumerate(things):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.65+row*1.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.10,1.0,0.95,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+1.15,y+0.15,1.75,0.40,cn,sz=14,b=True,c=c)
    tb(s,x+1.15,y+0.55,1.75,0.30,en,sz=10,c=GRAY)
    tb(s,x+1.15,y+0.80,1.75,0.30,"…还有更多!",sz=9,c=DARK)
# Bottom — interactive prompt
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.45),Inches(9.2),Inches(0.95))
prompt.fill.solid(); prompt.fill.fore_color.rgb=JOBS; prompt.line.color.rgb=IDEA; prompt.line.width=Pt(2.5)
tb(s,0.55,4.52,9.0,0.32,"🙋 你还能想到什么? 举手抢答!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,4.84,9.0,0.30,"💬 「iPhone 还可以 ___ !」",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.14,9.0,0.26,"What else can it do? Hands up — let's list more!",sz=9,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — 老师 iPhone 互动:\n• 老师真的拿出自己的手机 — 学生看到实物会兴奋\n• 一个一个问: 「这个你用过吗?」\n  - 打电话 / 拍照 / 听音乐 / 看视频 / 玩游戏 / 查地图\n• 抢答 — 学生还能想到什么 (闹钟! 翻译! 计算器! 量身高! ...)\n• 关键: 「一个小小的手机 — 这么多功能!」\n• 引出下一页: 「但你知道吗? 20 年前还没有 iPhone! 那时候怎么办?」")

# === JOBS INTRO 2: Before iPhone — needed many separate gadgets ===
s=ns(); bg(s,CREAM); hb(s,"📲 iPhone 发明之前  Before iPhone",JOBS)
tb(s,0.4,0.85,9.2,0.40,"2007 年之前 — 想做这些事,得带一大堆东西!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Before 2007 — you needed all these SEPARATE gadgets!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 5 old gadgets you used to need
gadgets=[
    ("☎️","老式电话","Phone","只能打电话",LAB),
    ("📷","相机","Camera","只能拍照",JOBS),
    ("🎧","MP3 / Walkman","MP3","只能听歌",PURPLE),
    ("📺","小电视","TV","只能看",MONEY),
    ("🗺️","纸地图","Paper map","容易迷路!",SKY),
]
for i,(em,cn,en,detail,c) in enumerate(gadgets):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(1.78),Inches(2.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,1.78,1.70,0.95,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.80,1.70,0.36,cn,sz=13,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.18,1.70,0.26,en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,1.70,0.65,detail,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Big insight at bottom
take=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.50),Inches(9.2),Inches(0.95))
take.fill.solid(); take.fill.fore_color.rgb=JOBS; take.line.color.rgb=IDEA; take.line.width=Pt(2.5)
tb(s,0.55,4.58,9.0,0.32,"😵 出门要带: 电话 + 相机 + MP3 + 地图 + 电视… 口袋都装不下!",sz=13,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,4.92,9.0,0.30,"💡 有一个人就想: 「能不能把这些都装进一个小盒子里?」",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.22,"One man thought: 'Can we put them ALL in one little box?'",sz=9,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — 以前的世界:\n• 老师演一下: 「出门要带电话、还要带相机、还要带 MP3、还要带地图……」 (假装口袋装不下)\n• 戏剧效果 — 让学生笑出来\n• 问: 「这多不方便! 你觉得哪个最麻烦?」\n• 关键句子: 「有一个人 — 他看到了这个『需要』, 然后想 — 能不能把这些都装进一个小盒子里?」\n• 留悬念: 「这个人是谁? — 翻下一页听故事!」")

# === JOBS INTRO 3: iPhone story (read aloud from image) ===
s=ns(); bg(s,CREAM); hb(s,"📖 乔布斯的故事  Steve Jobs' Story",JOBS)
tb(s,0.4,0.85,9.2,0.36,"听老师讲故事 — 看看乔布斯是怎么想出 iPhone 的!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Listen to the story — how did Jobs come up with the iPhone?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Embed the iPhone story image
img_path=os.path.join(os.path.dirname(__file__),"D3 resources","iPhone story.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path,Inches(0.6),Inches(1.55),width=Inches(8.8),height=Inches(3.40))
else:
    ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(1.55),Inches(8.8),Inches(3.40))
    ph.fill.solid(); ph.fill.fore_color.rgb=IMGBG; ph.line.color.rgb=JOBS; ph.line.width=Pt(2)
    tb(s,0.6,3.00,8.8,0.5,"📖 iPhone story.png (D3 resources/)",sz=18,b=True,c=LGRAY,a=PP_ALIGN.CENTER)
# Bottom transition strip
bridge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.00),Inches(9.2),Inches(0.50))
bridge.fill.solid(); bridge.fill.fore_color.rgb=BIZ; bridge.line.color.rgb=IDEA; bridge.line.width=Pt(2)
tb(s,0.55,5.05,9.0,0.30,"💼 乔布斯 = 企业家 — 看见需要 → 做出产品 → 帮助几亿人!",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.30,9.0,0.18,"Jobs = entrepreneur — saw needs → made products → helped billions!",sz=8,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5-6 分钟 — 讲 iPhone 故事:\n• 老师把 'iPhone story.png' 念给学生听\n• 念慢一点 — 让学生看图 + 听故事\n• 念完后问:\n  - 「乔布斯看到了什么需要?」 → 出门要带太多东西\n  - 「他做出了什么产品?」 → iPhone\n  - 「现在谁在用 iPhone?」 → 全世界几亿人\n• 关键 takeaway: 「乔布斯就是企业家 — 因为他看见需要、做出产品、帮助别人!」\n• 引出下一页: 「我们来讨论几个问题!」\n• 资源: D3 resources/iPhone story.png")

# === DISCUSSION HELPER (6 question cards in 2x3 grid) ===
def discuss_slide(title_cn,title_en,subtitle_cn,subtitle_en,questions,accent):
    """questions: list of 6 tuples (emoji, q_cn, q_en)"""
    s=ns(); bg(s,CREAM); hb(s,f"💬 讨论 · {title_cn}  ·  Discussion · {title_en}",accent)
    tb(s,0.4,0.85,9.2,0.32,subtitle_cn,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.18,9.2,0.22,subtitle_en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    # 2x3 grid of question cards (3 cols × 2 rows)
    for i,(em,q_cn,q_en) in enumerate(questions):
        col=i%3; row=i//3
        x=0.4+col*3.10; y=1.55+row*1.85
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.70))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=accent; sh.line.width=Pt(2.5)
        # Number badge
        nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.45),Inches(0.45))
        nb.fill.solid(); nb.fill.fore_color.rgb=accent; nb.line.fill.background()
        tb(s,x+0.10,y+0.14,0.45,0.36,str(i+1),sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        # Emoji
        tb(s,x+0.65,y+0.10,0.55,0.45,em,sz=24,a=PP_ALIGN.CENTER)
        # Question CN
        tb(s,x+0.15,y+0.65,2.65,0.65,q_cn,sz=12,b=True,c=DARK)
        # Question EN
        tb(s,x+0.15,y+1.32,2.65,0.32,q_en,sz=8,c=GRAY)
    # Bottom prompt
    bottom=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.25),Inches(9.2),Inches(0.30))
    bottom.fill.solid(); bottom.fill.fore_color.rgb=accent; bottom.line.fill.background()
    tb(s,0.5,5.29,9.0,0.22,"🙋 老师选 1-2 个问题,全班讨论或同桌交流!",sz=10,b=True,c=IDEA,a=PP_ALIGN.CENTER)
    return s

# === DISCUSSION 1: 关于乔布斯 ===
s=discuss_slide("关于乔布斯","About Jobs",
                "听完故事,一起来想想这些问题:",
                "Now that you've heard the story — discuss these questions:",
                [
                    ("🔍","乔布斯发现了什么问题?","What problem did Jobs discover?"),
                    ("📲","他为什么觉得带这么多东西很麻烦?","Why was carrying so many things annoying?"),
                    ("💡","他想到了什么解决办法?","What was his solution idea?"),
                    ("🚫","为什么很多人觉得他的想法「不可能」?","Why did people say his idea was 'impossible'?"),
                    ("📱","他想设计一部什么样的手机?","What kind of phone did Jobs want?"),
                    ("✨","为什么他说「简单,就是最厉害的设计」?","Why did he say 'simple is the best design'?"),
                ],JOBS)
n+=1; pn(s,n)
notes(s,"5-7 分钟 — 关于乔布斯的讨论:\n• 老师选 1-2 个问题全班讨论\n• 参考答案:\n  - 1 发现: 「出门要带电话、相机、MP3、地图…… 太麻烦了!」\n  - 2 麻烦: 「口袋装不下! 容易丢 / 忘带 / 充电也麻烦」\n  - 3 办法: 「把所有东西装进一个小盒子 — iPhone!」\n  - 4 不可能: 「以前没人做过 — 这么多功能怎么放一起?」\n  - 5 设计: 「漂亮 + 简单 + 只有 1 个按钮」\n  - 6 简单: 「太复杂 — 没人会用; 简单 — 大家都会用!」\n• 全班跟读关键句: 「简单,就是最厉害的设计!」")

# === DISCUSSION 2: 关于改进 + 不放弃 ===
s=discuss_slide("关于改进 + 不放弃","About Iteration",
                "iPhone 也不是一次就成功 — 看看他们是怎么一步步改进的:",
                "iPhone wasn't perfect at first — let's see how they improved:",
                [
                    ("❌","iPhone 第一次就成功了吗?","Did iPhone succeed the first time?"),
                    ("🔧","手机出问题的时候,他们怎么做?","What did they do when problems came up?"),
                    ("🔄","他们为什么要一直改、改了又改?","Why did they revise so many times?"),
                    ("💪","你觉得乔布斯放弃过吗? 为什么?","Did Jobs give up? Why or why not?"),
                    ("🌟","iPhone 帮人们做了哪些事?","How does iPhone help people today?"),
                    ("❤️","为什么这么多人喜欢 iPhone?","Why do so many people love iPhone?"),
                ],PURPLE)
n+=1; pn(s,n)
notes(s,"5-7 分钟 — 关于改进 + 不放弃的讨论:\n• 选 1-2 个问题讨论\n• 参考答案:\n  - 1 第一次: 不是! 早期 iPhone 有很多 bug, 信号不好, 电池不耐用\n  - 2 出问题: 工程师和设计师一起改进 — 找原因、想办法、测试\n  - 3 一直改: 「一次不完美 → 改一改 → 再试! 改了好几百次!」\n  - 4 放弃: 「没有! 他相信自己的想法 — 一直改到完美」\n  - 5 帮人: 打电话 / 拍照 / 找路 / 听歌 / 学习 / 看视频 / 联系家人\n  - 6 喜欢: 简单 + 漂亮 + 一个顶十个 (功能多但不复杂)\n• 关键 message: 「企业家 = 不放弃 + 一直改进!」")

# === DISCUSSION 3: 关于你 — APPLY TO YOURSELF ===
s=discuss_slide("关于你","About YOU",
                "现在轮到你想一想 — 你也能当小企业家!",
                "Now think — you can be a little entrepreneur too!",
                [
                    ("💼","听完故事,你觉得企业家在做什么?","From this story, what do entrepreneurs do?"),
                    ("⭐","你觉得企业家最重要的本领是什么?","What's the most important entrepreneur skill?"),
                    ("🤔","你发现了一个「小麻烦」,你会怎么解决?","If YOU saw a small problem, how would you fix it?"),
                    ("🎁","你想发明什么来帮助别人?","What would YOU invent to help others?"),
                    ("🔁","你觉得「失败了再试一次」重要吗?","Is 'try again after failing' important?"),
                    ("📱","如果你是设计师,你想给手机加什么新功能?","If YOU were a designer — what feature would you add to the phone?"),
                ],BIZ)
n+=1; pn(s,n)
notes(s,"6-8 分钟 — 把故事用到自己身上 (重点!):\n• 这是最重要的一组问题 — 让学生从故事 → 联系到自己\n• 选 2-3 个问题 (建议第 3、4、6 — 最能启发创意)\n• 同桌交流 (Turn & Talk): 2 人一组分享想法\n• 让 3-4 个学生上台分享 (每人 1 句)\n• 老师不评好坏 — 鼓励每一个想法\n• 关键句型:\n  - 「我发现 ___ 的问题, 我想 ___」\n  - 「我想发明 ___, 帮 ___」\n  - 「我想给手机加 ___ 功能」\n• 收集学生想法 — 可以留到 Session 3 设计产品时用!\n• 引出下一页: 「现在我们来看看 — 乔布斯是怎么一步一步做出 iPhone 的!」")

# === JOBS: 3 SIMPLE STEPS (K-5 friendly — mirrors slide 19's 3-card layout) ===
s=ns(); bg(s,CREAM); hb(s,"💡 乔布斯是怎么做的?  How Did Jobs Do It?",JOBS)
tb(s,0.4,0.85,9.2,0.40,"3 个简单的步骤 — 看见需要 → 做出来 → 大家爱用!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"3 simple steps — Saw a need → Made it → People loved it!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3 example cards (mirrors slide 19's layout)
jobs3=[
    ("👀","看见需要","Saw a Need","出门要带电话、相机、地图…… 太多了!",LAB),
    ("📱","做出 iPhone","Made iPhone","全部装进一个 — 又简单又漂亮!",JOBS),
    ("🎉","大家爱用","People Love It","全世界几亿人都在用!",MONEY),
]
for i,(em,cn,en,detail,c) in enumerate(jobs3):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.95),Inches(2.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,1.78,2.85,1.00,em,sz=62,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.85,2.85,0.40,cn,sz=18,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.25,2.85,0.28,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.60,2.85,0.55,detail,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Bottom takeaway strip
take=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.40),Inches(9.2),Inches(1.05))
take.fill.solid(); take.fill.fore_color.rgb=JOBS; take.line.color.rgb=IDEA; take.line.width=Pt(2.5)
tb(s,0.55,4.48,9.0,0.32,"🌟 这就是企业家 — 看见需要 → 做出来 → 帮助大家!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,4.82,9.0,0.26,"That's an entrepreneur — see a need → make it → help everyone!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.55,5.14,9.0,0.26,"💬 老师先问: 「他看到了什么需要?」 — 让学生猜!",sz=11,b=True,c=IDEA,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — Jobs 的 3 个简单步骤 (适合 K-5):\n• Step 1 (看见需要): 老师先问 —「他看到了什么需要?」 让学生举手猜!\n  - 提示: 「2007 年之前 — 你想打电话、拍照、听音乐、查地图…… 要带几样东西?」\n  - 学生猜完后揭晓: 「太多东西要带 — 又重又麻烦!」\n• Step 2 (做出 iPhone): 「他把所有功能装进一个小盒子 — iPhone! 又简单又漂亮!」\n• Step 3 (大家爱用): 「现在全世界几亿人都在用 — 他帮了几亿人!」\n• 关键 takeaway: 「企业家 = 看见 → 做 → 大家爱用」\n• 引出下一页: 「不只是乔布斯, 所有企业家都这样做!」")

# --- JOBS VIDEO INTRO (moved up: right after Jobs 3-step process) ---
s=ns(); bg(s,CREAM); hb(s,"📱 苹果之父 · 乔布斯  Father of Apple · Steve Jobs",JOBS)
tb(s,0.4,0.85,9.2,0.40,"乔布斯 (1955–2011) — 让电脑和手机变得又简单又好看!",sz=18,b=True,c=JOBS,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Jobs made computers + phones simple and beautiful — Apple & iPhone!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
vsh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.65),Inches(9.0),Inches(2.85))
vsh.fill.solid(); vsh.fill.fore_color.rgb=DARK; vsh.line.color.rgb=JOBS; vsh.line.width=Pt(3)
tb(s,0.5,2.05,9.0,1.40,"▶",sz=130,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.55,9.0,0.30,"📱 2007 年 — 第一部 iPhone 诞生!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.95,9.0,0.30,"老师在这里插入乔布斯视频 / Teacher: insert Jobs video",sz=10,b=True,c=IDEA,a=PP_ALIGN.CENTER)
hint=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
hint.fill.solid(); hint.fill.fore_color.rgb=WARM; hint.line.color.rgb=JOBS; hint.line.width=Pt(2)
tb(s,0.55,4.70,9.0,0.30,"🔍 视频建议: 'Steve Jobs for kids' / 'iPhone first launch' (2-3 分钟)",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,5.02,9.0,0.30,"👂 看的时候想一想: 乔布斯看到了什么「需要」?",sz=11,b=True,c=JOBS,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-5 分钟 — 介绍乔布斯:\n• 老师播放 Jobs 视频 (2-3 分钟) — 最好用 iPhone 发布会的片段\n• 推荐链接: youtube.com/watch?v=LOb3FJhDbYs (2007 iPhone Keynote)\n• 让学生看: 乔布斯看到了什么需要? 他做出了什么产品?\n• 看完串场: 「记得吗 — 我们刚学了乔布斯的 3 步 — 看见需要 → 做出来 → 大家爱用!」\n• 引出下一页 — 你来当乔布斯! 设计一个新功能!")

# ============================================================
# === SESSION 1 · PHASE 3: ACTIVE PRACTICE (25 min) ===
# Hands-on design + kid entrepreneurs
# ============================================================
# --- JOBS: HANDS-ON — design your imaginary phone feature ---
s=ns(); bg(s,CREAM); hb(s,"📱 你来当乔布斯!  Be Jobs: Design 1 New Feature",JOBS)
tb(s,0.4,0.85,9.2,0.40,"如果你是乔布斯 — 你想给手机加一个什么新功能?",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"If YOU were Jobs — what new feature would you add to a phone?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3 example feature cards (inspire)
features=[
    ("🐕","狗狗翻译器","Dog Translator","狗一叫 → 变成文字!",PURPLE),
    ("🧊","自动出冰水","Auto Cold","按一下 → 冰凉的水",SKY),
    ("👨‍👩‍👧","一键找家人","Find Family","按一下 → 知道家人在哪",MONEY),
]
for i,(em,cn,en,detail,c) in enumerate(features):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.95),Inches(2.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,1.78,2.85,0.85,em,sz=52,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.70,2.85,0.36,cn,sz=15,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.08,2.85,0.26,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.40,2.85,0.36,detail,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Activity prompt
act=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.95),Inches(9.2),Inches(0.95))
act.fill.solid(); act.fill.fore_color.rgb=JOBS; act.line.color.rgb=IDEA; act.line.width=Pt(2.5)
tb(s,0.55,4.02,9.0,0.30,"✏️ 3 分钟 — 想一个新功能,画下来 + 写: 谁会用? 解决什么问题?",sz=13,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,4.36,9.0,0.26,"3 min — invent 1 new feature, sketch it, name a customer + problem!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.55,4.65,9.0,0.28,"🌟 分享: 2-3 位学生上台 — 同学鼓掌 = 「顾客喜欢!」",sz=11,b=True,c=IDEA,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,5.00,"我设计了 ___ 功能 — 帮 ___ 解决 ___ 的问题。","I designed ___ feature — helps ___ solve ___.",accent=JOBS)
n+=1; pn(s,n)
notes(s,"5-6 分钟 — 你来当乔布斯:\n• 每个学生拿一张小纸\n• 1 分钟: 看 3 个例子, 找灵感\n• 3 分钟: 个人画 + 写 — 你想加什么新功能?\n  - 叫什么名字? (狗狗翻译器、一键找妈妈……)\n  - 谁会用? (狗主人、小朋友、老人……)\n  - 解决什么问题?\n• 2 分钟: 2-3 个学生上台分享 — 用句型说\n• 同学鼓掌 = 「顾客喜欢!」 — 掌声越响, 想用的人越多")

# --- KID ENTREPRENEURS GALLERY (2 cards — Lily Born + 9-yr popcorn CEO) ---
s=ns(); bg(s,CREAM); hb(s,"🧒 两位和你一样的小老板  Kid Bosses Like YOU",PH_CLOSE)
tb(s,0.4,0.85,9.2,0.34,"这两位小朋友 — 看见身边的「需要」,就当上老板了!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.24,"These 2 kids saw needs near them — and became bosses!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 2 kid entrepreneur cards (side by side, large)
kids=[
    ("🦘","Lily Born","8 岁开始 · 美国","三脚杯 (Kangaroo Cup)","爷爷手抖,水总是洒出来!","Grandpa's hands shake — spills!","Amazon 上卖了 2 万多个!",LAB),
    ("🍿","9 岁男孩","9 岁 · 美国","爆米花公司 (Popcorn CEO!)","小朋友想吃更好吃的爆米花","Kids want yummy popcorn!","卖到全美国!",MAYUN),
]
for i,(em,name,age,product,need_cn,need_en,success,c) in enumerate(kids):
    x=0.4+i*4.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(4.55),Inches(3.05))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.10,1.68,4.35,1.0,em,sz=72,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,2.75,4.35,0.36,name,sz=18,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.12,4.35,0.26,age,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.42,4.35,0.32,product,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.78,4.35,0.30,need_cn,sz=11,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.05,4.35,0.20,need_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    # success pill
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.30),Inches(4.28),Inches(3.95),Inches(0.28))
    sp.fill.solid(); sp.fill.fore_color.rgb=IDEA; sp.line.fill.background()
    tb(s,x+0.30,4.31,3.95,0.22,f"🏆 {success}",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Interactive prompt bar
vote=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.75),Inches(9.2),Inches(0.80))
vote.fill.solid(); vote.fill.fore_color.rgb=BIZ; vote.line.color.rgb=IDEA; vote.line.width=Pt(2.5)
tb(s,0.55,4.82,9.0,0.30,"🌟 「他们也是小朋友 — 你也可以!」 你想像谁一样? 解决什么问题?",sz=13,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.15,9.0,0.30,"💬 「我想像 ___ 一样, 解决 ___ 的需要!」",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5-7 分钟 — 两位小朋友企业家:\n• 🦘 Lily Born (8 岁开始, 现在 15 岁): 爷爷有帕金森, 喝水老是洒 → 设计了「三脚杯」(Kangaroo Cup, 不会倒的杯子)! 现在在 Amazon 卖了 2 万多个!\n  - 资源: seinsights.asia/article/5504\n  - 关键故事: 在自己家里测试 prototype, 全家人帮她提建议\n• 🍿 9 岁男孩爆米花 CEO: 美国一个 9 岁男孩自己开爆米花公司, 卖到全美国\n  - 资源: youtube.com/watch?v=pVx3AUOHY0c (下一页视频)\n\n• 互动 (3 分钟):\n  - 让 2-3 个学生说: 「我想像 ___ 一样, 解决 ___」\n  - 提示: 你身边有没有像「爷爷手抖」这样的真问题?\n\n• 关键 message:\n  - 「他们也是小朋友 — 跟你们一样!」\n  - 「他们解决的都是身边的小事 (爷爷喝水、朋友想吃的零食)」\n  - 「你也可以 — 一会儿我们就来试!」\n• 接下来: 翻到下一页, 听听 Lily 的杯子故事 — 她是怎么帮爷爷的!")

# --- LILY'S CUP STORY (dedicated deep-dive — read aloud from image) ---
s=ns(); bg(s,CREAM); hb(s,"📖 Lily 的杯子故事  Lily's Cup Story",LAB)
tb(s,0.4,0.85,9.2,0.36,"听老师讲故事 — 8 岁的 Lily 是怎么帮爷爷的?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Listen — how did 8-year-old Lily help her grandpa?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Embed the cup story image
img_path=os.path.join(os.path.dirname(__file__),"D3 resources","three leg cup story.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path,Inches(0.6),Inches(1.55),width=Inches(8.8),height=Inches(3.40))
else:
    ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(1.55),Inches(8.8),Inches(3.40))
    ph.fill.solid(); ph.fill.fore_color.rgb=IMGBG; ph.line.color.rgb=LAB; ph.line.width=Pt(2)
    tb(s,0.6,3.00,8.8,0.5,"📖 three leg cup story.png (D3 resources/)",sz=18,b=True,c=LGRAY,a=PP_ALIGN.CENTER)
# Bottom transition strip
bridge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.00),Inches(9.2),Inches(0.50))
bridge.fill.solid(); bridge.fill.fore_color.rgb=LAB; bridge.line.color.rgb=IDEA; bridge.line.width=Pt(2)
tb(s,0.55,5.05,9.0,0.30,"🦘 Lily = 小企业家 — 看见爷爷的需要 → 设计三脚杯 → 帮助手抖的人!",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.30,9.0,0.18,"Lily = entrepreneur — saw grandpa's need → made Kangaroo Cup → helps people with shaky hands!",sz=8,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5-6 分钟 — 讲 Lily 的杯子故事:\n• 老师把 'three leg cup story.png' 念给学生听\n• 念慢一点 — 让学生看图 + 听故事\n• 念完后问 4 个问题:\n  - 「Lily 看到了什么需要?」 → 爷爷有帕金森, 手抖, 喝水老洒出来\n  - 「她那时候几岁?」 → 8 岁! 跟你们差不多大!\n  - 「她做出了什么产品?」 → Kangaroo Cup (三脚杯, 三只脚, 不会倒)\n  - 「她现在卖了多少个?」 → Amazon 上 2 万多个!\n• 关键 takeaway: 「Lily 也是企业家 — 看见需要 (爷爷手抖) → 做出产品 (三脚杯) → 帮助别人 (所有手抖的人)!」\n• 关键 message: 「她 8 岁就开始了 — 你现在也可以!」\n• 全班跟读: 「我也能像 Lily 一样, 看见需要 → 做出东西 → 帮助别人!」\n• 资源: D3 resources/three leg cup story.png\n• 引出下一页: 「再来看一个 — 9 岁男孩爆米花 CEO 的视频!」")

# --- KID ENTREPRENEURS VIDEO INTRO (relatable!) ---
s=ns(); bg(s,CREAM); hb(s,"🌟 小朋友也能当企业家!  Kids Are Entrepreneurs Too!",PH_CLOSE)
tb(s,0.4,0.85,9.2,0.40,"再看一个 — 9 岁男孩,自己开公司卖爆米花!",sz=18,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"One more — a 9-year-old boy running his own popcorn company!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
vsh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.65),Inches(9.0),Inches(2.85))
vsh.fill.solid(); vsh.fill.fore_color.rgb=DARK; vsh.line.color.rgb=PH_CLOSE; vsh.line.width=Pt(3)
tb(s,0.5,2.05,9.0,1.40,"▶",sz=130,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.55,9.0,0.30,"🍿 美国 9 岁男孩 · 爆米花公司小老板!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.95,9.0,0.30,"老师在这里插入视频 / Teacher: insert kid entrepreneur video",sz=10,b=True,c=IDEA,a=PP_ALIGN.CENTER)
hint=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
hint.fill.solid(); hint.fill.fore_color.rgb=WARM; hint.line.color.rgb=PH_CLOSE; hint.line.width=Pt(2)
tb(s,0.55,4.70,9.0,0.30,"🔍 视频: youtube.com/watch?v=pVx3AUOHY0c (9 岁男孩爆米花 CEO · 2-3 分钟)",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,5.02,9.0,0.30,"👂 看的时候想想: 他看到了什么需要? 他几岁开始的?",sz=11,b=True,c=BIZ,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-5 分钟 — 9 岁男孩爆米花 CEO 视频:\n• 老师播放 YouTube 视频 (链接在 slide 上)\n• 让学生看后说: 「这个小朋友几岁? 他卖什么?」\n• 串场: 「Lily 8 岁 + 这个男孩 9 岁 — 都比你们大不了多少!」\n• 关键 message: 「当企业家不用等长大 — 你现在就可以!」")

# --- SHARED TRAITS SUMMARY ---
s=ns(); bg(s,CREAM); hb(s,"🌟 企业家的共同点  Entrepreneur Traits",NAVY)
tb(s,0.4,0.85,9.2,0.40,"乔布斯、Lily、爆米花 CEO — 年龄不同, 但他们都……",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Jobs, Lily, the popcorn boy — different ages, but they ALL share these:",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
traits=[
    ("👀","看见需要","See Needs","别人没看到的,他们看到了",NAVY),
    ("💡","想点子","Got Ideas","脑筋一直在转",IDEA),
    ("🔄","不放弃","Don't Give Up","失败 → 改 → 再试",LAB),
    ("🤝","帮助顾客","Help Customers","让别人喜欢用",MONEY),
    ("💼","勇敢去试","Take Risks","敢做新东西",BIZ),
]
for i,(em,cn,en,detail,c) in enumerate(traits):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(1.78),Inches(2.95))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,1.78,1.70,0.85,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,1.70,0.40,cn,sz=14,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,1.70,0.30,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.50,1.70,0.85,detail,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
take=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.75),Inches(9.2),Inches(0.85))
take.fill.solid(); take.fill.fore_color.rgb=NAVY; take.line.color.rgb=IDEA; take.line.width=Pt(2.5)
tb(s,0.55,4.82,9.0,0.32,"🌟 这些特点你也有 — 你也能当小企业家!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.18,9.0,0.30,"You have these traits too — YOU can be a little entrepreneur!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — 共同点小结:\n• 一个一个念 5 个 traits, 让学生跟读\n• 让 1-2 个学生说: 「我有 ___ 这个特点!」\n• 关键: 「这些都是性格 — 不是天生的, 是练出来的!」\n• 「这节课你已经用过了 — 看 4 个生活问题 + 设计新功能!」\n• 引出下一节: 「现在轮到你试一试 — 你看到了什么需要?」")

# ============================================================
# === SESSION 1 · PHASE 4: APPLY (10 min) ===
# ============================================================
# Apply 1 — 我看到的"需要" (whole class brainstorm with 4 categories)
s=ns(); bg(s,CREAM); hb(s,"💭 我看到的「需要」  Needs I See",GREEN)
tb(s,0.4,0.85,9.2,0.36,"你在哪里看到过「需要」? 看看下面这 4 个地方的例子!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Where have YOU seen needs? See 4 places of examples!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
need_categories=[
    ("🏫","学校 School",["铅笔总丢","排队太吵","午饭排队太慢","书包太重"],GREEN),
    ("🏠","家 Home",["弟弟妹妹太闹","妈妈太累","遥控器找不到","房间太乱"],BIZ),
    ("🌍","世界 World",["塑料垃圾太多","老人没人陪","流浪动物没家","下雨没伞"],SKY),
    ("👫","朋友 Friends",["朋友不会中文","同学没吃早餐","新同学没朋友","作业不会做"],PURPLE),
]
for i,(em,cn_label,examples,cl) in enumerate(need_categories):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,x+0.10,y+0.10,0.7,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.13,3.6,0.40,cn_label,sz=14,b=True,c=cl)
    for j,ex in enumerate(examples):
        tb(s,x+0.20,y+0.62+j*0.22,4.20,0.22,f"·  {ex}",sz=10,c=DARK)
disc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.70))
disc.fill.solid(); disc.fill.fore_color.rgb=GREEN; disc.line.fill.background()
tb(s,0.55,4.92,9.0,0.30,"👥 3-4 人一组讨论 — 每人说一个自己见过的需要!",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.30,"💬 我看到 ___ 的需要 — 有一个问题是 ___ 。",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🌱 APPLY 1 · 5 分钟 — 分组头脑风暴:\n• Step 1 (1 分钟): 个人看例子, 想一个自己见过的需要\n• Step 2 (3 分钟): 3-4 人一组, 每人说一个需要\n  - 句型: 「我看到 ___ 的需要 — 有一个问题是 ___」\n• Step 3 (1 分钟): 每组选 1 个「最想解决的」, 准备下一页想办法\n• 不评判好坏 — 鼓励每个想法\n• 提示: 选的需要 = 一会儿小生意的灵感!")

# Apply 2 — 我的小生意 idea
s=ns(); bg(s,CREAM); hb(s,"💡 我的小生意  My Mini-Business Idea",BIZ)
tb(s,0.4,0.85,9.2,0.36,"用你刚选的需要 — 想一个小生意 (产品或服务)!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Use the need you picked — invent a mini-business (product OR service)!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 4 fields to fill in
fields=[
    ("📛","名字","Name","「我的生意叫 ___ 」",BIZ),
    ("👥","顾客","Customer","「___ 会用我的东西」",PURPLE),
    ("🎯","解决","Problem","「解决 ___ 的问题」",LAB),
    ("💰","为什么","Why Win","「因为我的 ___ 」",MONEY),
]
for i,(em,cn,en,frame,c) in enumerate(fields):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.15,0.7,0.50,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.10,3.6,0.40,cn,sz=15,b=True,c=c)
    tb(s,x+0.85,y+0.45,3.6,0.28,en,sz=9,c=GRAY)
    tb(s,x+0.15,y+0.82,4.30,0.50,frame,sz=12,b=True,c=DARK)
disc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.70))
disc.fill.solid(); disc.fill.fore_color.rgb=BIZ; disc.line.fill.background()
tb(s,0.55,4.92,9.0,0.30,"✏️ 个人想 5 分钟,然后跟同桌分享!",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.30,"5 min individual brainstorm → share with 1 classmate!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🌱 APPLY 2 · 5-6 分钟 — 个人 brainstorm:\n• 1 分钟 — 老师演示 1 个例子:\n  - 名字:「找鞋小狗」  顾客:「家里的妹妹」  解决:「鞋老是找不到」  为什么:「它闻一下就知道」\n• 4 分钟 — 学生个人填 4 个 fields (画也行)\n• 1-2 分钟 — 跟同桌分享 (Turn & Talk)\n• 重要: 不评好坏! 这只是先想一想 — Session 3 才完整设计")

# ============================================================
# === SESSION 1 · PHASE 5: SHARE & CLOSE (5 min) ===
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎤 总结 + 下一步  Summary + Next!",PH_CLOSE)
score_badge(s)
tb(s,0.4,0.85,9.2,0.4,"🧭 今天早上学了什么?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.40),Inches(4.4),Inches(1.85))
left.fill.solid(); left.fill.fore_color.rgb=BIZ; left.line.fill.background()
tb(s,0.5,1.50,4.4,0.5,"💡 企业家是谁?",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.00,4.4,0.55,"看见需要!",sz=22,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,2.55,4.4,0.30,"See needs · 乔布斯 + Lily + 爆米花 CEO",sz=11,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.85,4.4,0.30,"做出产品 + 帮助顾客",sz=11,c=IDEA,a=PP_ALIGN.CENTER)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.40),Inches(4.4),Inches(1.85))
right.fill.solid(); right.fill.fore_color.rgb=MONEY; right.line.fill.background()
tb(s,5.10,1.50,4.4,0.5,"🌟 共同点",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,5.10,2.00,4.4,0.55,"不放弃!",sz=22,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,5.10,2.55,4.4,0.30,"Don't give up · 失败 → 改 → 再试",sz=11,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,5.10,2.85,4.4,0.30,"勇敢去试 · 帮助别人",sz=11,c=IDEA,a=PP_ALIGN.CENTER)
trans=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.40),Inches(9.0),Inches(1.45))
trans.fill.solid(); trans.fill.fore_color.rgb=PH_ACTIVE; trans.line.color.rgb=IDEA; trans.line.width=Pt(3)
tb(s,0.5,3.50,9.0,0.45,"🛠️ 下午的项目课 → 你来当老板!",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.95,9.0,0.30,"Afternoon Project Class → YOU become the boss!",sz=11,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,4.30,9.0,0.40,"💼 4 个项目: 发现 → 设计 → 推销 → 帮学校发明!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎤 SHARE & CLOSE · 5 分钟:\n• 1 分钟 — 全班齐喊: 「企业家! 看见需要! 做出产品! 帮助别人!」 (3 遍, 加手势)\n• 1 分钟 — 每队代表说一句: 「我们队学到了 ___ 」\n• 1 分钟 — 公布 Session 1 暂时积分\n• 2 分钟 — Tease 下午: 「下午 4 个项目: 发现需要、设计产品、卖给同学, 还要给学校发明一个帮助工具!」")

# ============================================================
# 22. SESSION 2 DIVIDER
# ============================================================
s=div("Session 2  下午 1:00–1:50","📚 复习 + 我会认 + 我会写  Review · Read · Write",IDEA,"📖"); n+=1; pn(s,n)

# ============================================================
# 23. REVIEW — Session 1 recap
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔄 复习  Review · Session 1",NAVY)
tb(s,0.4,0.85,9.2,0.40,"早上 学了 什么？  What did we learn this morning?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Quick recap before we read & write!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
recap=[
    ("💡","企业家 = 看见需要","Entrepreneur sees needs","→ 做出产品 + 帮助顾客",BIZ),
    ("📱","乔布斯 · iPhone","Jobs · iPhone","以前带一大堆 → 现在一个 iPhone",JOBS),
    ("🦘","Lily · 三脚杯","Lily · Kangaroo Cup","爷爷手抖 → 设计不会倒的杯子",LAB),
    ("🍿","9 岁爆米花 CEO","9-yr Popcorn CEO","小朋友也能当老板!",MAYUN),
]
for i,(em,cn,en,detail,c) in enumerate(recap):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.65+row*1.45
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.30))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.12,0.7,0.65,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.10,3.6,0.40,cn,sz=15,b=True,c=c)
    tb(s,x+0.85,y+0.45,3.6,0.30,en,sz=10,c=GRAY)
    tb(s,x+0.15,y+0.85,4.30,0.40,detail,sz=11,b=True,c=DARK)
sentence_frame_bar(s,4.65,
    "今天早上我学了 ___ 。",
    "This morning I learned about ___.")
n+=1; pn(s,n)
notes(s,"5 分钟 review:\n• 「早上学了几位企业家? 都是谁?」 抢答\n• 复习: 企业家 = 看见需要 → 做出产品 → 帮助顾客\n• 让 1-2 个学生说: 「我今天学了 ___」")

# ============================================================
# 24-29. 我会认 — 6 vocabulary words: 企业家、产品、顾客、买、卖、钱
# ============================================================
read_words=[
    ("💼","企业家","qǐ yè jiā","Entrepreneur",BIZ,
        "乔布斯是一位很厉害的企业家。",
        "📷 企业家 / 西装 / 公司"),
    ("📦","产品","chǎn pǐn","Product",JOBS,
        "iPhone 是苹果公司的产品。",
        "📷 玩具 / 手机 / 食物"),
    ("👥","顾客","gù kè","Customer",PURPLE,
        "顾客喜欢用,生意就成功啦!",
        "📷 顾客 / 收银台 / 笑脸"),
    ("🛒","买","mǎi","Buy",MONEY,
        "我用钱买了一个玩具。",
        "📷 购物车 / 商场 / 钱"),
    ("🏪","卖","mài","Sell",LEGO,
        "妈妈把蛋糕卖给顾客。",
        "📷 摊位 / 店 / 收钱"),
    ("💰","钱","qián","Money",IDEA,
        "顾客给老板钱,老板给顾客产品。",
        "📷 硬币 / 纸币 / 钱包"),
]
for em,cn,py,en,c,sent,img_label in read_words:
    s=ns(); bg(s,CREAM); hb(s,f"👀 我会认 · {cn}  I Can Read",c)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.fill.background()
    tb(s,0.5,1.10,4.3,1.4,cn,sz=60 if len(cn)==3 else 70,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.40,4.3,0.4,f"{py}  {en}",sz=18,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读!  Read after me!",sz=14,c=c,a=PP_ALIGN.CENTER)
    ib_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.3),Inches(1.0),Inches(4.4),Inches(2.5))
    ib_box.fill.solid(); ib_box.fill.fore_color.rgb=IMGBG; ib_box.line.fill.background()
    tb(s,5.3,2.05,4.4,0.4,img_label,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=c; sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=c)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    n+=1; pn(s,n)
    notes(s,f"3 分钟 — {cn}:\n• 老师指字, 全班齐读 3 遍 (慢→快→唱)\n• 看图: 「这是 ___, 在做什么?」\n• 读例句, 学生跟读\n• 抽 1-2 个学生用「{cn}」造一个新句子\n• 写到黑板上 — 让学生在空中跟着写一遍")

# ============================================================
# 30. 我会写 · 产品 (stroke order)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 产品  I Can Write · Product",BIZ)
tb(s,0.4,0.85,9.2,0.40,"一起来写「产品」!",sz=22,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 产品 (Product)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tianzi(s,0.55,1.65,2.20,"产",BIZ,pinyin="chǎn (produce)",char_sz=120)
tianzi(s,2.95,1.65,2.20,"品",BIZ,pinyin="pǐn (item)",char_sz=120)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
panel.line.color.rgb=BIZ; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=BIZ; head.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写  How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"1️⃣「产」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  上面是一点+丿, 下面是「厂」",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"2️⃣「品」 — 9 笔 (3 个口!)",sz=14,b=True,c=DARK)
tb(s,5.45,3.40,4.10,0.30,"  3 个「口」 — 上面 1 个,下面 2 个",sz=10,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格里写 3 遍",sz=12,b=True,c=BIZ)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「产品」! 我的产品叫 ___ 。",
    "I can write 产品! My product is called ___.")
n+=1; pn(s,n)
notes(s,"5-6 分钟:\n• 演示笔顺, 学生跟写 (空中写)\n• 田字格练 3 遍\n• 记忆法: 「品」 = 3 个「口」 — 上 1 下 2, 像金字塔!\n• 让学生说说自己想设计的产品 — 「我的产品叫 ___ 」")

# ============================================================
# 31. 我会写 · 买 (stroke order)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 买  I Can Write · Buy",MONEY)
tb(s,0.4,0.85,9.2,0.40,"一起来写「买」!",sz=22,b=True,c=MONEY,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 买 (Buy)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tianzi(s,1.30,1.65,2.95,"买",MONEY,pinyin="mǎi (buy)",char_sz=160)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
panel.line.color.rgb=MONEY; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=MONEY; head.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写  How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"📐「买」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  上面是「乛+丶」, 下面是「头」",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"💡 记忆: 「头」+ 上面一个帽子 = 「买」",sz=12,b=True,c=MONEY)
tb(s,5.45,3.40,4.10,0.30,"  Like 「头」 with a hat on top!",sz=9,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格里写 3 遍",sz=12,b=True,c=MONEY)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「买」! 我想买 ___ 。",
    "I can write 买! I want to buy ___.",accent=MONEY)
n+=1; pn(s,n)
notes(s,"5 分钟:\n• 演示笔顺\n• 记忆法: 「买」像「头」上戴了顶帽子 — 戴上帽子出门去买东西!\n• 注意跟「卖」的区别 (下一页)\n• 让学生说: 「我想买 ___」")

# ============================================================
# 32. 我会写 · 卖 (stroke order)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 卖  I Can Write · Sell",LEGO)
tb(s,0.4,0.85,9.2,0.40,"一起来写「卖」!",sz=22,b=True,c=LEGO,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 卖 (Sell) — 比「买」多一个「十」!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tianzi(s,1.30,1.65,2.95,"卖",LEGO,pinyin="mài (sell)",char_sz=160)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
panel.line.color.rgb=LEGO; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=LEGO; head.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写  How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"📐「卖」 — 8 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  上面「十」+「买」 = 卖!",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"💡 记忆: 卖 = 买 + 十",sz=12,b=True,c=LEGO)
tb(s,5.45,3.40,4.10,0.30,"  卖东西比买东西多个「+10」 (赚钱了!)",sz=9,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格里写 3 遍",sz=12,b=True,c=LEGO)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「卖」! 我想卖 ___ 。",
    "I can write 卖! I want to sell ___.",accent=LEGO)
n+=1; pn(s,n)
notes(s,"5-6 分钟:\n• 演示笔顺 — 强调比「买」多一个「十」在上面\n• 记忆法: 「卖 = 买 + 十 — 卖东西比买东西多 (赚到钱了!)」\n• 全班一起空写「买」「卖」对比\n• 让学生说: 「我想卖 ___」 — 引到 Session 3 的小生意\n• 写完 → 「我会写」贴纸; 写得最好的队 +2 分")

# ============================================================
# 33. SESSION 3 DIVIDER — Mini-Business Projects
# ============================================================
s=div("Session 3  下午 2:00–2:50","🛠️ 项目课  迷你小生意!  ·  50 min",PH_ACTIVE,"💼"); n+=1; pn(s,n)

# ============================================================
# === SESSION 3 · PHASE 1: HOOK (5 min) ===
# ============================================================
# Hook — classroom problem demo (parallel to Day 2's paper bridge hook)
s=ns(); bg(s,CREAM); hb(s,"🤔 教室里的一个真问题  A Real Classroom Problem",LAB)
tb(s,0.4,0.85,9.2,0.40,"老师现在带来一个真实的问题 — 你来想想怎么解决?",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Teacher shows a real classroom problem — how would YOU solve it?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Problem display (left)
prob=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.70),Inches(4.50),Inches(2.85))
prob.fill.solid(); prob.fill.fore_color.rgb=LAB; prob.line.color.rgb=IDEA; prob.line.width=Pt(3)
tb(s,0.5,1.90,4.50,1.0,"📚",sz=80,a=PP_ALIGN.CENTER)
tb(s,0.5,2.85,4.50,0.40,"教室里的书太多 — 没地方放!",sz=16,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.25,4.50,0.30,"Too many books — no place to put them!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.5,3.65,4.50,0.40,"❓ 你有什么好办法?",sz=15,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,4.05,4.50,0.30,"What ideas do YOU have?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
# Idea bubble (right)
idea_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.15),Inches(1.70),Inches(4.40),Inches(2.85))
idea_box.fill.solid(); idea_box.fill.fore_color.rgb=WHITE; idea_box.line.color.rgb=BIZ; idea_box.line.width=Pt(3)
tb(s,5.30,1.80,4.10,0.40,"💡 头脑风暴",sz=15,b=True,c=BIZ)
tb(s,5.30,2.15,4.10,0.30,"全班一起想 30 秒!",sz=11,c=GRAY)
ideas=["📦 做一个漂亮的书箱?",
       "🔄 开一个借书服务?",
       "📅 每周换一批书?",
       "🏪 开个二手书小铺?"]
for i,it in enumerate(ideas):
    tb(s,5.30,2.55+i*0.42,4.10,0.35,it,sz=12,b=True,c=DARK)
tpr_strip(s,4.65,"举手说: 我想 ___ !","Raise hand: I want to ___!")
n+=1; pn(s,n)
notes(s,"🔥 HOOK · 5 分钟:\n• 老师演一下: 抱一大堆书走过来 — 「啊! 书太多了! 没地方放!」\n• 戏剧效果\n• 问: 「你来帮老师想想办法?」\n• 全班 30 秒头脑风暴 — 收集 4-5 个想法\n• 提示例子 (学生卡住时): 做书箱? 借书服务? 二手书铺?\n• 引出: 「这就是企业家 — 看见一个问题, 想一个办法! 接下来 4 个项目, 你自己来试!」")

# ============================================================
# === SESSION 3 · PHASE 2: MINI-LESSON (10 min) ===
# ============================================================
# Mini 1 — 企业家 4 step loop (action version for Session 3)
s=ns(); bg(s,CREAM); hb(s,"🔁 企业家 4 步循环  Entrepreneur Loop",BIZ)
score_badge(s)
tb(s,0.4,0.85,9.2,0.4,"看见 → 想 → 做 → 试! 失败了就 改 → 再试!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.30,"See → Think → Make → Test! Fail → Improve → Try again!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("👀","看","See","看见什么需要?",NAVY),
    ("💡","想","Think","怎么帮?",BIZ),
    ("🛠️","做","Make","画出来 + 做出来!",PH_ACTIVE),
    ("🤝","试","Test","顾客喜不喜欢?",MONEY),
]
for i,(em,cn,en,q,cl) in enumerate(steps):
    x=0.4+i*2.30
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.10),Inches(2.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.78),Inches(1.85),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+0.78,1.93,0.55,0.4,str(i+1),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.50,2.0,0.6,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.20,2.0,0.45,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.0,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.95,2.0,0.30,q,sz=11,c=DARK,a=PP_ALIGN.CENTER)
    if i<3:
        tb(s,x+2.10,2.85,0.30,0.4,"→",sz=22,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,4.40,9.2,0.30,"🔁 顾客不喜欢? 改一改 → 再试! (就像乔布斯一样!)",sz=14,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,4.70,9.2,0.25,"Customer doesn't like it? Improve → try again! (just like Jobs!)",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"📚 MINI 1 · 4 分钟:\n• 4 步走一遍 — 强调第 4 步「试」 → 「改」最重要\n• 「乔布斯改了好几百次!」\n• 全班齐声 + 手势: 看! 想! 做! 试! (看=指眼睛, 想=指头, 做=拳头, 试=举手)")

# Mini 2 — 4 projects overview
s=ns(); bg(s,CREAM); hb(s,"💼 今天的 4 个项目  4 Projects Today",BIZ)
tb(s,0.4,0.85,9.2,0.4,"全班一起做 4 个项目 — 一个接一个!",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"4 mini-projects today — one after another!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
projects=[
    ("P1","🔍","发现需要","Find a Need","5 min","头脑风暴 — 你看到了什么需要?",BIZ),
    ("P2","✏️","设计产品","Design Product","8 min","画一画 + 写下来: 名字、顾客、价钱!",JOBS),
    ("P3","🛒","小小推销员","Mini Sales Fair","7 min","上台介绍 + 同学投票!",MONEY),
    ("P4","🟣","帮助学校","Help School","10 min","小组设计: 解决一个学校问题!",PURPLE),
]
for i,(tag,em,cn,en,t,desc,cl) in enumerate(projects):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.22),Inches(3.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=cl; sh.line.fill.background()
    tb(s,x+0.05,1.85,2.12,0.40,tag,sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.25,2.12,0.80,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.10,2.12,0.36,cn,sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.46,2.12,0.26,en,sz=9,c=IDEA,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.78,2.12,0.30,f"⏱  {t}",sz=11,b=True,c=IDEA,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,4.12,2.02,1.0,desc,sz=9,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"📚 MINI 2 · 2 分钟:\n• 介绍 4 个项目 — 流程: P1 → P2 → P3 → P4\n• 强调总时间 30 分钟 active practice + apply\n• P1+P2+P3 = 个人/小组的「小生意」 — 用早上学的 4 步\n• P4 = 大家一起解决一个学校的真问题\n• 每完成一个项目 = 队 +5 分")

# ============================================================
# === SESSION 3 · PHASE 3: ACTIVE PRACTICE (20 min) ===
# Projects 1-3
# ============================================================
# Project 1 — 发现 需要 (brainstorm)
s=ns(); bg(s,CREAM); hb(s,"🔍 项目 1 · 发现需要  Project 1 · Find a Need",BIZ)
group_label(s)
score_badge(s)
tb(s,0.4,1.20,9.2,0.40,"⏱ 5 分钟 — 每队选出一个「最想解决」的需要!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Need source quadrants
sources=[
    ("🏫","学校","School","铅笔丢 / 排队吵 / 排队吃饭",GREEN),
    ("🏠","家","Home","遥控器丢 / 玩具乱 / 妈妈累",BIZ),
    ("👫","同学","Friends","新同学没朋友 / 作业难",PURPLE),
    ("🌍","世界","World","垃圾多 / 老人孤单 / 没伞",SKY),
]
for i,(em,cn,en,detail,c) in enumerate(sources):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.70+row*1.40
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.30))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.15,1.0,0.95,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+1.20,y+0.10,3.30,0.42,cn,sz=16,b=True,c=c)
    tb(s,x+1.20,y+0.52,3.30,0.28,en,sz=10,c=GRAY)
    tb(s,x+1.20,y+0.82,3.30,0.40,detail,sz=11,b=True,c=DARK)
sentence_frame_bar(s,4.60,"我们队选了 ___ 这个需要 — 因为 ___ 。","Our team picked ___ — because ___.",accent=BIZ)
n+=1; pn(s,n)
notes(s,"🎯 P1 · 5 分钟:\n• 分 4 队, 每队 4-5 人\n• 1 分钟 — 每人想一个需要 (从 4 个来源选)\n• 3 分钟 — 队内投票, 选 1 个「最想解决的」\n• 1 分钟 — 队代表上台用句型报告\n• 提示: 选的这个需要就是 P2 设计产品的起点!\n• 每队 +1 分 (参与)")

# Project 2 — 设计 产品 (design card)
s=ns(); bg(s,CREAM); hb(s,"✏️ 项目 2 · 设计产品  Project 2 · Design Product",JOBS)
group_label(s)
score_badge(s)
tb(s,0.4,1.20,9.2,0.40,"⏱ 8 分钟 — 画 + 写出你的「产品名片」!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Product card template (4 quadrants)
card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.70),Inches(9.2),Inches(3.10))
card.fill.solid(); card.fill.fore_color.rgb=WARM; card.line.color.rgb=JOBS; card.line.width=Pt(3)
tb(s,0.4,1.78,9.2,0.36,"📋 产品名片模板  Product Card Template",sz=15,b=True,c=JOBS,a=PP_ALIGN.CENTER)
fields=[
    ("📛","名字 + 图","Name + Drawing","「我的产品叫 ___」",JOBS),
    ("👥","顾客是谁","Who's the Customer","「___ 会用它」",PURPLE),
    ("🎯","解决什么","Problem","「解决 ___ 的问题」",LAB),
    ("💰","卖多少钱","Price","「卖 ___ 元」",MONEY),
]
for i,(em,cn,en,frame,c) in enumerate(fields):
    col=i%2; row=i//2
    x=0.55+col*4.45; y=2.20+row*1.20
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.35),Inches(1.05))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.15,0.55,0.5,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.70,y+0.08,3.6,0.32,cn,sz=12,b=True,c=c)
    tb(s,x+0.70,y+0.38,3.6,0.24,en,sz=8,c=GRAY)
    tb(s,x+0.15,y+0.68,4.10,0.36,frame,sz=10,b=True,c=DARK)
sentence_frame_bar(s,4.85,"我的产品是 ___, 卖给 ___, 帮他们 ___, 卖 ___ 元!","My product is ___, for ___, helps ___, sells for ___!",accent=JOBS)
n+=1; pn(s,n)
notes(s,"🎯 P2 · 8 分钟:\n• 老师发 A4 纸 (或「产品名片」模板)\n• 4 分钟 — 队内设计:\n  - 画产品 (中央一个大图)\n  - 写名字 (上方)\n  - 写顾客 + 解决什么 (左 + 右)\n  - 定价钱 (1 元 / 5 元 / 10 元)\n• 4 分钟 — 队内练推销词: 「我的产品是 ___, 卖给 ___, 帮他们 ___, 卖 ___ 元」\n• 每队 +2 分 (完成产品名片)")

# Project 3 — 小小 推销员 (mini fair)
s=ns(); bg(s,CREAM); hb(s,"🛒 项目 3 · 小小推销员  Project 3 · Mini Sales Fair",MONEY)
group_label(s)
score_badge(s)
tb(s,0.4,1.20,9.2,0.40,"⏱ 7 分钟 — 各队上台推销 + 同学投票!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
# 3 step columns
steps=[
    ("1","🎤","上台推销","Pitch on stage","用句型介绍 60 秒",MONEY),
    ("2","👀","同学看听","Audience listens","顾客可以提 1 个问题",PURPLE),
    ("3","🗳️","顾客投票","Customers vote","用假钱或贴纸投给最想买的!",BIZ),
]
for i,(num,em,cn,en,detail,c) in enumerate(steps):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(2.95),Inches(2.80))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    pill(s,x+1.18,1.80,0.60,0.36,num,c,sz=13)
    tb(s,x+0.05,2.25,2.85,0.70,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.00,2.85,0.36,cn,sz=15,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.36,2.85,0.26,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,3.70,2.80,0.7,detail,sz=10,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Vote bar
vote=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.60),Inches(9.2),Inches(0.85))
vote.fill.solid(); vote.fill.fore_color.rgb=MONEY; vote.line.color.rgb=IDEA; vote.line.width=Pt(2.5)
tb(s,0.55,4.68,9.0,0.30,"🏆 票数最多的队 = 「最受欢迎产品」 + 5 分!",sz=13,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,5.02,9.0,0.30,"Most votes = 'Most Popular Product' + 5 pts!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎯 P3 · 7 分钟 — Mini Fair:\n• 准备假钱 / 投票贴纸 — 每个学生发 2 张\n• 每队 60 秒上台推销 (4 队 × 60 = 4 分钟)\n• 1 分钟 — 顾客 (其他队) 走一圈, 把「钱」/ 贴纸投给自己最想买的\n• 2 分钟 — 数票, 公布「最受欢迎产品」\n• 各队都有奖: 最受欢迎 / 最有创意 / 最实用 / 最团结 → 每队 +5 分!\n• 重要: 这不是比输赢 — 我们要庆祝每一个想法!")

# ============================================================
# === SESSION 3 · PHASE 4: APPLY (10 min) ===
# Project 4: 发明 帮助 学校 的 东西
# ============================================================
# Project 4 — 发明 帮助 学校 的 东西
s=ns(); bg(s,CREAM); hb(s,"🟣 项目 4 · 发明工具帮学校!  Project 4 · Help Our School",PURPLE)
group_label(s)
score_badge(s)
tb(s,0.4,1.20,9.2,0.36,"先讨论: 我们学校里有什么问题? 然后设计解决方案!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.55,9.2,0.26,"Step 1: discuss real school problems · Step 2: design a solution!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Left: 4 school problems (the user's specific examples)
problems_left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.95),Inches(4.55),Inches(2.85))
problems_left.fill.solid(); problems_left.fill.fore_color.rgb=WARM; problems_left.line.color.rgb=PURPLE; problems_left.line.width=Pt(3)
tb(s,0.55,2.05,4.30,0.36,"🏫 学校里的真问题",sz=14,b=True,c=PURPLE)
tb(s,0.55,2.40,4.30,0.24,"Real problems at school",sz=9,c=GRAY)
school_probs=[("✏️","铅笔总丢","Pencils always lost"),
              ("🔊","排队太吵","Lining up too noisy"),
              ("🎒","书包太乱","Backpack messy"),
              ("🍱","午饭排队太慢","Lunch line too slow")]
for i,(em,cn,en) in enumerate(school_probs):
    y=2.70+i*0.50
    tb(s,0.70,y,0.50,0.40,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.02,3.60,0.28,cn,sz=12,b=True,c=DARK)
    tb(s,1.30,y+0.28,3.60,0.20,en,sz=8,c=GRAY)
# Right: 4 steps to invent solution
sol=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.95),Inches(4.55),Inches(2.85))
sol.fill.solid(); sol.fill.fore_color.rgb=WHITE; sol.line.color.rgb=PURPLE; sol.line.width=Pt(3)
tb(s,5.20,2.05,4.30,0.36,"💡 4 步设计解决方案",sz=14,b=True,c=PURPLE)
tb(s,5.20,2.40,4.30,0.24,"4 steps to invent a solution",sz=9,c=GRAY)
sol_steps=[("1️⃣","选一个问题","Pick 1 problem"),
           ("2️⃣","想一想怎么解决","Think — how to fix?"),
           ("3️⃣","画下来 + 起名字","Draw + name it"),
           ("4️⃣","队代表上台分享!","Team rep shares!")]
for i,(num,cn,en) in enumerate(sol_steps):
    y=2.70+i*0.50
    tb(s,5.20,y,0.55,0.40,num,sz=14,b=True,c=PURPLE)
    tb(s,5.80,y+0.02,3.65,0.28,cn,sz=12,b=True,c=DARK)
    tb(s,5.80,y+0.28,3.65,0.20,en,sz=8,c=GRAY)
sentence_frame_bar(s,4.85,"我们队发明了 ___ — 解决 ___ 的问题, 帮 ___ 。","Our team invents ___ — solves ___, helps ___.",accent=PURPLE)
n+=1; pn(s,n)
notes(s,"🟣 P4 · 10 分钟 — 帮学校发明:\n• 这是最重要的项目 — 真正用上企业家的思维!\n• Step 1 (2 分钟): 队内讨论 — 学校里还有哪些真问题?\n  - 不限于 PPT 上 4 个 — 鼓励学生自己说真见过的\n• Step 2 (3 分钟): 每队选 1 个问题 + 头脑风暴解决方案\n  - 老师例子: 铅笔丢 → 「铅笔姓名贴 + 班级失物角」\n  - 排队吵 → 「排队拍节奏游戏」\n  - 书包乱 → 「分格书包 / 整理小袋子」\n  - 午饭慢 → 「分时段 / 提前点餐」\n• Step 3 (3 分钟): 画出解决方案 + 起名字\n• Step 4 (2 分钟): 队代表上台分享\n• 重点强调: 这个项目「特别像企业家」 — 因为企业家就是在解决问题!\n• 老师可以挑一个想法真的给老师 / 校长 — 让学生觉得自己的想法被重视!\n• 每队 +5 分 (完成)")

# ============================================================
# === SESSION 3 · PHASE 5: SHARE & CLOSE (5 min) ===
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🏆 颁奖 + Day 3 徽章!  Awards + Day 3 Badge!",PH_CLOSE)
score_badge(s)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(1.45))
sh.fill.solid(); sh.fill.fore_color.rgb=PH_CLOSE; sh.line.color.rgb=DARK; sh.line.width=Pt(2.5)
tb(s,0.4,1.00,9.2,0.45,"🏆 今天的总冠军队  Day Champions",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.50,9.2,0.45,"___ 队 — 全天最高分!",sz=24,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,0.4,2.00,9.2,0.30,"___ team — highest score across all 3 sessions!",sz=11,c=DARK,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.55),Inches(3),Inches(2.6))
sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=IDEA; sh.line.width=Pt(6)
tf=tb(s,3.6,2.75,2.8,2.4,"DAY 3",sz=18,b=True,c=IDEA,a=PP_ALIGN.CENTER)
ap(tf,"💼💡",sz=42,a=PP_ALIGN.CENTER)
ap(tf,"小小企业家",sz=14,b=True,c=BIZ,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=12,b=True,c=OK,a=PP_ALIGN.CENTER)
ap(tf,"📱🦘🍿",sz=18,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎤 CLOSE · 5 分钟:\n• 1 分钟 — 公布全天总分, 冠军队上台\n• 1 分钟 — 各队代表说一句: 「我们今天最喜欢 ___」\n• 2 分钟 — 给每个学生发徽章 / 贴纸 — 「你也是小企业家!」\n• 1 分钟 — 全班齐喊: 「看见需要! 做出产品! 帮助别人! 我也是小企业家!」\n• 预告 Day 4: 「明天 — Community Helpers!」")

# ============================================================
# BONUS SUPPLEMENTARY VIDEO (end-of-day inspiration)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎬 加餐视频  Bonus Video — More Inspiration!",PH_CLOSE)
tb(s,0.4,0.85,9.2,0.40,"还有时间? 我们再看一个视频, 收获更多灵感!",sz=16,b=True,c=BIZ,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"If time permits — one more video for extra inspiration!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
vsh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.65),Inches(9.0),Inches(2.85))
vsh.fill.solid(); vsh.fill.fore_color.rgb=DARK; vsh.line.color.rgb=PH_CLOSE; vsh.line.width=Pt(3)
tb(s,0.5,2.05,9.0,1.40,"▶",sz=130,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.55,9.0,0.30,"🌟 更多企业家 / 小老板的故事 — 给学生更多启发!",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.95,9.0,0.30,"老师在这里插入视频 / Teacher: insert supplementary video",sz=10,b=True,c=IDEA,a=PP_ALIGN.CENTER)
hint=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.85))
hint.fill.solid(); hint.fill.fore_color.rgb=WARM; hint.line.color.rgb=PH_CLOSE; hint.line.width=Pt(2)
tb(s,0.55,4.70,9.0,0.30,"🔍 视频: youtube.com/watch?v=MnrJzXM7a6o (补充视频 / supplementary)",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,5.02,9.0,0.30,"💭 看完想一想: 你今天学到的哪一点 — 最让你想去试?",sz=11,b=True,c=BIZ,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎬 加餐视频 · 3-5 分钟 (可选):\n• 如果还有时间 (Session 3 提前结束 / 课与课之间的 buffer), 播放此视频\n• 链接: youtube.com/watch?v=MnrJzXM7a6o\n• 看完让 1-2 个学生分享: 「我今天学到 ___, 我想试 ___」\n• 这是一整天的收尾 — 让学生带着「我也可以」的心情回家!")

# ============================================================
out=os.path.join(os.path.dirname(__file__),"day3_entrepreneurs.pptx")
prs.save(out); print(f"Saved {out}  ({n} slides)")
