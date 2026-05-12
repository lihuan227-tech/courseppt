#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
我的职业梦想 — Day 4: 社区小帮手 (Community Helper Experience Day)
2 sessions × 50 min · K-5 混龄, ~20 学生 · 中文沉浸式 · 高互动 / role play / project-based

Session 1 (50 min) — 🎭 Community Helper Experience Day
  · Warm-up 猜职业动作 (5)
  · 绘本引入 (10)
  · Role Play 一日小老师 / 一日小医生 / 一日小厨师 (30)
  · 总结 + Reflection (5)

Session 2 (50 min) — 复习 + 语言目标 + Projects
  · 复习 / 完成 (15-20)
  · 我会认 老师·学校·医院·消防员·厨师 + 我会写 老师·学校 (15)
  · Project 1 社区地图 (小组) + Project 2 Thank You 奖状 (个人) + 颁奖典礼 (15-20)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# === Palette ===
HELP   = RGBColor(0x2E,0x7D,0x32)   # community green
HEART  = RGBColor(0xE5,0x3E,0x5E)   # warm pink-red
NAVY   = RGBColor(0x1E,0x3A,0x5F)
IDEA   = RGBColor(0xF5,0xC2,0x42)
LAB    = RGBColor(0xE5,0x3E,0x3E)
GREEN  = RGBColor(0x2E,0x7D,0x32)
PURPLE = RGBColor(0x7B,0x1F,0xA2)
RUST   = RGBColor(0xC4,0x52,0x2A)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
BROWN  = RGBColor(0x6B,0x44,0x23)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
OK     = RGBColor(0x38,0x8E,0x3C)
SKY    = RGBColor(0x1F,0x77,0xB4)
GOLD   = RGBColor(0xF5,0xA6,0x23)

# Phase colors
PH_WARM   = GOLD
PH_BOOK   = SKY
PH_PLAY   = RUST
PH_LANG   = PURPLE
PH_PROJ   = HELP
PH_CLOSE  = IDEA

# Per-helper colors
TEACH  = RGBColor(0x43,0xA0,0x47)   # teacher green
DOC    = RGBColor(0xC8,0x25,0x3E)   # doctor wine red
CHEF   = RGBColor(0xFF,0x8F,0x00)   # chef amber
FIRE   = RGBColor(0xD3,0x18,0x18)   # firefighter red

# === Helpers ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])

def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b
    r.font.color.rgb=c; r.font.name='KaiTi'
    return tf

def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b
    r.font.color.rgb=c; r.font.name='KaiTi'

def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)

def hb(s,txt,c=HELP,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)

def pn(s,n):
    tb(s,9.0,5.28,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)

def time_pill(s,minutes,t=0.20,l=8.20):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(1.55),Inches(0.36))
    sh.fill.solid(); sh.fill.fore_color.rgb=IDEA; sh.line.fill.background()
    tb(s,l,t+0.04,1.55,0.32,f"⏱  {minutes} 分钟",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)

def notes(s,text):
    nf=s.notes_slide.notes_text_frame; lines=text.split("\n")
    nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line

def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    return s

def sentence_frame_bar(s,t,frame_cn,frame_en,accent=IDEA):
    if t > 4.95: t = 4.95
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid(); sf.fill.fore_color.rgb=WARM; sf.line.color.rgb=accent; sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.7,0.4,"💬 我来说",sz=14,b=True,c=accent)
    tb(s,2.0,t+0.07,7.6,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.0,t+0.32,7.6,0.3,frame_en,sz=10,c=GRAY)

def phase_marker(emoji,phase_cn,phase_en,time_min,color,what_cn,what_en):
    s=ns(); bg(s,color)
    tb(s,1,0.85,8,0.7,emoji,sz=80,a=PP_ALIGN.CENTER)
    tb(s,1,1.85,8,0.6,phase_cn,sz=38,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.50,8,0.4,phase_en,sz=16,c=IDEA,a=PP_ALIGN.CENTER)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.5),Inches(3.10),Inches(3.0),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=IDEA; sh.line.fill.background()
    tb(s,3.5,3.18,3.0,0.4,f"⏱  {time_min} 分钟",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,1,4.00,8,0.5,what_cn,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,4.55,8,0.35,what_en,sz=12,c=IDEA,a=PP_ALIGN.CENTER)
    return s

def emoji_circle(s,l,t,size,emoji,fill,border=IDEA,em_sz=88):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l),Inches(t),Inches(size),Inches(size))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.color.rgb=border; sh.line.width=Pt(4)
    tb(s,l,t+(size-em_sz/72)/2,size,size,emoji,sz=em_sz,a=PP_ALIGN.CENTER)

def tianzi(s,x,y,size,char,color,pinyin=None,char_sz=130):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(size),Inches(size))
    box.fill.solid(); box.fill.fore_color.rgb=WHITE
    box.line.color.rgb=color; box.line.width=Pt(3)
    mid_x=x+size/2; mid_y=y+size/2; lw=0.015
    v=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(mid_x-lw/2),Inches(y),Inches(lw),Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb=LGRAY; v.line.fill.background()
    h=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(mid_y-lw/2),Inches(size),Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb=LGRAY; h.line.fill.background()
    tb(s,x,y,size,size,char,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s,x,y+size+0.05,size,0.30,pinyin,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)

# Three-panel activity layout: 老师引导语 · 学生任务 · 讨论问题
def activity_panels(s,top,teacher_cn,student_cn,discuss_cn,accent=HELP):
    panels=[("👩‍🏫","老师引导语",teacher_cn,accent),
            ("👧","学生任务",student_cn,RUST),
            ("💭","讨论",discuss_cn,SKY)]
    for i,(em,label,text,cl) in enumerate(panels):
        x=0.40+i*3.10
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(top),Inches(2.95),Inches(1.85))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
        head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(top),Inches(2.95),Inches(0.42))
        head.fill.solid(); head.fill.fore_color.rgb=cl; head.line.fill.background()
        tb(s,x+0.10,top+0.06,2.75,0.32,f"{em} {label}",sz=12,b=True,c=WHITE)
        tb(s,x+0.12,top+0.48,2.71,1.30,text,sz=11,b=True,c=DARK)

# Scenario role-play card
def scenario_card(s,top,scenario_em,scenario_cn,scenario_en,pro_action,sentence_frames,color):
    # Left: big scenario
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(top),Inches(3.90),Inches(3.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.color.rgb=IDEA; sh.line.width=Pt(3)
    tb(s,0.40,top+0.20,3.90,1.20,scenario_em,sz=86,a=PP_ALIGN.CENTER)
    tb(s,0.40,top+1.60,3.90,0.50,scenario_cn,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.40,top+2.10,3.90,0.32,scenario_en,sz=11,c=IDEA,a=PP_ALIGN.CENTER)
    tb(s,0.40,top+2.55,3.90,0.34,"❓ 你怎么办?",sz=15,b=True,c=IDEA,a=PP_ALIGN.CENTER)
    tb(s,0.40,top+2.90,3.90,0.30,"What would YOU do?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
    # Right top: Professional action box
    rt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.45),Inches(top),Inches(5.15),Inches(1.55))
    rt.fill.solid(); rt.fill.fore_color.rgb=WHITE; rt.line.color.rgb=color; rt.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.45),Inches(top),Inches(5.15),Inches(0.42))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,4.55,top+0.06,5.0,0.32,"🌟 Professional 会怎么做",sz=12,b=True,c=WHITE)
    tb(s,4.60,top+0.48,4.95,1.05,pro_action,sz=12,b=True,c=DARK)
    # Right bottom: Sentence frames
    rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.45),Inches(top+1.65),Inches(5.15),Inches(1.75))
    rb.fill.solid(); rb.fill.fore_color.rgb=WARM; rb.line.color.rgb=color; rb.line.width=Pt(2.5)
    head2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.45),Inches(top+1.65),Inches(5.15),Inches(0.42))
    head2.fill.solid(); head2.fill.fore_color.rgb=color; head2.line.fill.background()
    tb(s,4.55,top+1.71,5.0,0.32,"💬 可用句型",sz=12,b=True,c=WHITE)
    for i,frame in enumerate(sentence_frames):
        tb(s,4.65,top+2.13+i*0.32,4.90,0.30,f"• {frame}",sz=12,b=True,c=DARK)

n=0

# ============================================================
# 1. COVER
# ============================================================
s=ns(); bg(s,HELP)
tb(s,1,0.50,8,0.50,"我的职业梦想 · My Dream Career",sz=18,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,1,1.05,8,0.85,"Day 4",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.50,8,0.90,"🏘  社区小帮手",sz=54,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,2.60,8,0.50,"Community Helper Experience Day",sz=20,b=True,c=IDEA,a=PP_ALIGN.CENTER)
# 4 helper emoji circles
helpers=[("👩‍🏫",TEACH,3.50),("👩‍⚕️",DOC,4.65),("👨‍🍳",CHEF,5.80),("🚒",FIRE,6.95)]
for em,cl,xpos in helpers:
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(xpos-1.10),Inches(3.35),Inches(1.10),Inches(1.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=cl; sh.line.color.rgb=IDEA; sh.line.width=Pt(3)
    tb(s,xpos-1.10,3.45,1.10,1.0,em,sz=48,a=PP_ALIGN.CENTER)
tb(s,1,4.60,8,0.40,"🎭 体验 · 角色扮演 · 项目实践",sz=16,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,1,5.05,8,0.30,"K-5 混龄 · 约 20 人 · 2 sessions × 50 min",sz=11,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"Day 4 开场 (1 分钟):\n• 「今天 — 我们不只听故事, 我们要 变成 社区小帮手!」\n• 介绍 4 位主角: 老师 · 医生 · 厨师 · 消防员\n• 全班按 4 队分好 (红/蓝/绿/黄) — 准备 role play")

# ============================================================
# 2. 5-DAY PREVIEW (Day 4 highlighted)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🗺️ 5 天的职业之旅  Our 5-Day Career Journey",NAVY)
tb(s,0.4,0.85,9.2,0.34,"今天是第 4 天 — Community Helper Experience Day!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.28,"Day 4 — meet, role-play, and become helpers!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
days_preview=[
    ("Day 1","认识职业世界","Discover Careers","🌍",NAVY,"8 个职业"),
    ("Day 2","小小科学家","Little Scientists","🔬",SKY,"⭐ 爱迪生"),
    ("Day 3","小小企业家","Little Entrepreneurs","💡",GOLD,"⭐ 乔布斯"),
    ("Day 4","社区小帮手","Community Helpers","🏘",HELP,"⭐ 今天!"),
    ("Day 5","AI 与未来","AI & the Future","🤖",PURPLE,"⭐ AI 公司"),
]
for i,(label,cn,en,em,cl,spotlight) in enumerate(days_preview):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.82),Inches(3.45))
    is_today=(i==3)
    sh.fill.solid(); sh.fill.fore_color.rgb=cl if is_today else WHITE
    sh.line.color.rgb=cl; sh.line.width=Pt(3.5 if is_today else 2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.1),Inches(1.65),Inches(0.55),Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb=WHITE if is_today else cl; badge.line.fill.background()
    tb(s,x+0.1,1.74,0.55,0.4,str(i+1),sz=18,b=True,c=cl if is_today else WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.7,1.7,1.1,0.3,label,sz=11,b=True,c=WHITE if is_today else cl)
    tb(s,x+0.05,2.30,1.72,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,1.72,0.4,cn,sz=15,b=True,c=WHITE if is_today else DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,1.72,0.3,en,sz=10,c=IDEA if is_today else GRAY,a=PP_ALIGN.CENTER)
    sep=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+0.25),Inches(3.85),Inches(1.32),Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb=WHITE if is_today else cl; sep.line.fill.background()
    tb(s,x+0.05,4.00,1.72,0.85,spotlight,sz=11,b=True,c=WHITE if is_today else cl,a=PP_ALIGN.CENTER)
tb(s,0.4,5.20,9.2,0.30,"📍 今天我们 体验 + 演 + 做! Today we experience + act + make!",sz=12,b=True,c=HELP,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30 秒过渡:\n• 提醒整个 unit 的位置\n• 「今天 — 第 4 天! 我们 体验 4 位 community helper, 然后 做 社区地图 + 颁奖!」")

# ============================================================
# 3. SESSION 1 DIVIDER
# ============================================================
s=div("Session 1  ·  50 min","🎭 Community Helper Experience Day  ·  体验 · 演 · 反思",HELP,"🏘"); n+=1; pn(s,n)

# ============================================================
# === SESSION 1 · PART 1: WARM-UP (5 min) ===
# ============================================================
s=phase_marker("🎬","Part 1 · Warm-up","猜职业动作游戏",5,PH_WARM,"老师演动作 → 学生猜职业","Teacher mimes → students guess")
n+=1; pn(s,n)

# 4-1. 猜职业动作游戏
s=ns(); bg(s,CREAM); hb(s,"🎭 猜职业动作!  Guess the Job!",PH_WARM)
time_pill(s,5)
tb(s,0.4,0.85,9.2,0.36,"老师 演 一个 动作 — 你 猜 是 什么 职业?",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.24,"Teacher mimes an action — what job is it?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 4 mime cards in 2x2 grid
mimes=[
    ("✍️","写 黑板","Writing on board","→ 老师 Teacher",TEACH),
    ("🩺","听 心跳","Listening to heart","→ 医生 Doctor",DOC),
    ("🍳","做 饭","Cooking","→ 厨师 Chef",CHEF),
    ("🔥","救 火","Putting out fire","→ 消防员 Firefighter",FIRE),
]
for i,(em,action_cn,action_en,answer,cl) in enumerate(mimes):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.62))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.12,1.30,1.30,em,sz=64,a=PP_ALIGN.CENTER)
    tb(s,x+1.55,y+0.18,2.90,0.40,action_cn,sz=18,b=True,c=cl)
    tb(s,x+1.55,y+0.60,2.90,0.30,action_en,sz=10,c=GRAY)
    tb(s,x+1.55,y+0.95,2.90,0.40,answer,sz=15,b=True,c=DARK)
tb(s,0.4,5.10,9.2,0.36,"🙋 看见动作 → 大声 喊 出 职业 名字! ⏱ 每个动作 1 分钟",sz=12,b=True,c=HELP,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎬 WARM-UP · 5 分钟 — 猜职业动作:\n• 老师戏剧化地演 4 个动作 (一个一个来):\n  - ✍️ 写黑板 → 假装拿粉笔在空中写\n  - 🩺 听心跳 → 假装戴听诊器\n  - 🍳 做饭 → 假装炒菜\n  - 🔥 救火 → 假装拿水管\n• 学生猜 — 大声喊出来 — 答对鼓掌!\n• 不要 给提示 — 让 学生 自己看 动作\n• Bonus: 抽 1-2 个学生 上来 演 — 全班 猜!\n• 关键 message: 「这 4 个人 — 都在 帮助 大家! 今天 我们 都来 体验!」\n• 课堂管理: 学生 起立 围着 老师 半圆形 — 互动 更好")

# 4-2. Turn & Talk
s=ns(); bg(s,CREAM); hb(s,"🗣️ Turn & Talk · 你 最想 做 什么 工作?",PH_WARM)
time_pill(s,2)
tb(s,0.4,0.85,9.2,0.50,"❓ 你 长大后 最想 做 什么 工作?",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.40,9.2,0.30,"What job do you most want to do when you grow up?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Big T&T instruction box
box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(1.95),Inches(7.0),Inches(2.30))
box.fill.solid(); box.fill.fore_color.rgb=WHITE; box.line.color.rgb=PH_WARM; box.line.width=Pt(3)
tb(s,1.5,2.10,7.0,0.50,"👥",sz=42,a=PP_ALIGN.CENTER)
tb(s,1.5,2.70,7.0,0.45,"找 一个 同桌  Find a partner",sz=20,b=True,c=PH_WARM,a=PP_ALIGN.CENTER)
tb(s,1.5,3.20,7.0,0.36,"⏱  每人 30 秒 · 轮流 说",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1.5,3.60,7.0,0.30,"30 sec each, take turns",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.55,"我 长大 想 当 ___, 因为 ___ 。","I want to be a ___ when I grow up, because ___.",accent=PH_WARM)
n+=1; pn(s,n)
notes(s,"🗣️ TURN & TALK · 2 分钟:\n• 找 同桌 (或左右邻居)\n• 老师 计时 — 每人 30 秒\n• 不评判 — 让每个孩子 都说\n• 老师 走动 — 听 1-2 对孩子 的对话\n• 结束后 抽 2 位 上台 分享\n• 关键 message: 「不管 你 想 当 什么 — 都 可以 帮助 别人!」")

# ============================================================
# === SESSION 1 · PART 2: 绘本引入 (10 min) ===
# ============================================================
s=phase_marker("📖","Part 2 · 绘本引入","大人 上班 都在 做 什么?",10,PH_BOOK,"问题 → 视频 → 观察 → 讨论","Question → Video → Observe → Discuss")
n+=1; pn(s,n)

# 5-1. 绘本前问题
s=ns(); bg(s,CREAM); hb(s,"❓ 大人 上班 都 在 做 什么?",PH_BOOK)
time_pill(s,2)
tb(s,0.4,0.85,9.2,0.36,"你 知道 爸爸 妈妈 / 大人 每天 去 上班 — 都 在 做 什么 吗?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Do you know what grown-ups do at work every day?",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 7 example bubbles
examples=[
    ("📞","打电话"),("💻","用 电脑"),("👥","开 会"),
    ("📚","教 学生"),("🩺","看 病人"),("🍳","做 饭"),("🚗","开 车"),
]
for i,(em,cn) in enumerate(examples):
    col=i%4; row=i//4
    x=0.4+col*2.32; y=1.65+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(1.30))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=PH_BOOK; sh.line.width=Pt(2)
    tb(s,x+0.05,y+0.10,2.10,0.75,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.85,2.10,0.40,cn,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,5.05,"我 爸爸/妈妈 上班 时 ___ 。","My parent ___ at work.",accent=PH_BOOK)
n+=1; pn(s,n)
notes(s,"❓ 绘本前 · 2 分钟:\n• 老师 真诚 地 问: 「你 知道 你 爸妈 每天 上班 做 什么 吗?」\n• 让 4-5 个 学生 分享 (举手)\n• 不评判 — 收 集 想法 (打电话/开会/用电脑/教书/做饭...)\n• 老师 总结: 「不同的 大人 — 不同的 工作!」\n• 引出 关键词 — 「这些 帮 社区 的 人, 我们 叫 — Community Helpers / 社区小帮手!」")

# 5-2. Community Helpers REVEAL
s=ns(); bg(s,CREAM); hb(s,"🌟 他们 叫 — Community Helpers!  社区小帮手",HELP)
tb(s,0.4,0.85,9.2,0.40,"很多 工作 — 都 在 帮助 别人!",sz=18,b=True,c=HELP,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Many jobs help other people!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Big definition banner
banner=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.75),Inches(9.0),Inches(2.50))
banner.fill.solid(); banner.fill.fore_color.rgb=HELP; banner.line.color.rgb=IDEA; banner.line.width=Pt(3)
tb(s,0.5,1.85,9.0,0.65,"Community Helpers",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.55,9.0,0.50,"社区 小 帮手",sz=28,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.15,9.0,0.40,"= 帮助 社区 的 人",sz=18,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.60,9.0,0.30,"People who help our community.",sz=11,c=WARM,a=PP_ALIGN.CENTER)
# Bottom transition
tr=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.45),Inches(9.2),Inches(0.90))
tr.fill.solid(); tr.fill.fore_color.rgb=WARM; tr.line.color.rgb=HELP; tr.line.width=Pt(2)
tb(s,0.55,4.52,9.0,0.30,"📖 现在 — 我们 一起 看 一个 视频 + 听 一个 故事!",sz=13,b=True,c=HELP,a=PP_ALIGN.CENTER)
tb(s,0.55,4.85,9.0,0.30,"Now — let's watch a video & hear a story!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.55,5.15,9.0,0.26,"听 的时候 — 想 一想: 有 哪些 community helpers?",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🌟 Reveal · 1 分钟:\n• 关键词 揭晓 — 全班 跟读 「社区小帮手 = Community Helpers」 3 遍\n• 关键 message: 「帮助 社区 的 人 — 就是 community helper」\n• 引出 视频")

# 5-3. Video — 大人上班都在做什么 (link)
s=ns(); bg(s,CREAM); hb(s,"🎬 看 视频  大人 上班 都在 做 什么?",PH_BOOK)
time_pill(s,4)
tb(s,0.4,0.85,9.2,0.38,"老师 播放 视频 — 一起 看!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"Teacher plays video — watch together!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
vsh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.60),Inches(9.0),Inches(2.85))
vsh.fill.solid(); vsh.fill.fore_color.rgb=DARK; vsh.line.color.rgb=PH_BOOK; vsh.line.width=Pt(3)
tb(s,0.5,2.00,9.0,1.40,"▶",sz=130,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.50,9.0,0.30,"📖 大人 上班 都在 做 什么?",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.85,9.0,0.30,"youtube.com/watch?v=7ARFgslUeJA",sz=11,b=True,c=IDEA,a=PP_ALIGN.CENTER)
link=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.55),Inches(9.2),Inches(0.85))
link.fill.solid(); link.fill.fore_color.rgb=WARM; link.line.color.rgb=PH_BOOK; link.line.width=Pt(2)
tb(s,0.55,4.62,9.0,0.30,"👀 看 的 时候 — 想 一想: 书里 出现了 哪些 工作?",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.30,"While watching — notice: what jobs do you see?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎬 视频 · 3-4 分钟:\n• 老师 提前 测试 视频 链接 + 投影 + 音量\n• 视频: youtube.com/watch?v=7ARFgslUeJA\n• 看 视频 前 强调: 「看 的 时候 — 数 一下 有 几个 工作!」\n• 视频 结束 — 不评论 — 直接 进入 下一页 观察 checklist")

# 5-4. 观察 checklist
s=ns(); bg(s,CREAM); hb(s,"👀 听 绘本 时 — 注意 看!",PH_BOOK)
time_pill(s,1)
tb(s,0.4,0.85,9.2,0.36,"看 视频 的 时候, 注意 这 3 件 事:",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"While watching, pay attention to 3 things:",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
checks=[
    ("1","📋","书里 出现 了 哪些 工作?","What jobs appear?",NAVY),
    ("2","🤝","他们 在 帮助 谁?","Who are they helping?",HELP),
    ("3","💡","他们 是 怎么 帮助 别人 的?","How do they help?",GOLD),
]
for i,(num,em,cn,en,cl) in enumerate(checks):
    y=1.65+i*1.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9.0),Inches(0.92))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.65),Inches(y+0.16),Inches(0.60),Inches(0.60))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.65,y+0.22,0.60,0.45,num,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.10,1.0,0.75,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,2.55,y+0.18,7.0,0.36,cn,sz=18,b=True,c=DARK)
    tb(s,2.55,y+0.58,7.0,0.26,en,sz=10,c=GRAY)
n+=1; pn(s,n)
notes(s,"👀 Observer checklist · 1 分钟:\n• 视频 前 — 老师 念 这 3 个问题\n• 让学生 心里 记住 — 「看 + 想」\n• 不写 — 只是 留意\n• 看 完后 — 用 下一页 讨论 问题")

# 5-5. 听完讨论问题
s=ns(); bg(s,CREAM); hb(s,"💭 听 完后 — 一起 讨论!  Discussion Questions",PH_BOOK)
time_pill(s,4)
tb(s,0.4,0.85,9.2,0.34,"选 1-2 个 问题 — 全班 讨论 / Think-Pair-Share!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.22,"Pick 1-2 questions — discuss as a class!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
qs=[
    ("🔍","你 看到了 哪些 工作?","What jobs did you see?"),
    ("⭐","哪个 工作 你 最 熟悉?","Which job is most familiar?"),
    ("📚","老师 是 怎么 帮助 学生 的?","How does a teacher help?"),
    ("🩺","医生 是 怎么 帮助 病人 的?","How does a doctor help?"),
    ("🍳","厨师 是 怎么 帮助 大家 的?","How does a chef help?"),
    ("😱","如果 没有 这些人 — 我们 的 生活 会 怎么样?","What if these helpers were gone?"),
]
for i,(em,cn,en) in enumerate(qs):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.55+row*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.70))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=PH_BOOK; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.45),Inches(0.45))
    nb.fill.solid(); nb.fill.fore_color.rgb=PH_BOOK; nb.line.fill.background()
    tb(s,x+0.10,y+0.14,0.45,0.36,str(i+1),sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.65,y+0.10,0.55,0.45,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,y+0.65,2.65,0.65,cn,sz=12,b=True,c=DARK)
    tb(s,x+0.15,y+1.32,2.65,0.32,en,sz=8,c=GRAY)
n+=1; pn(s,n)
notes(s,"💭 讨论 · 4 分钟:\n• 选 1-2 个 (建议 Q3-Q5 + Q6)\n• 用 Think-Pair-Share:\n  - 30 秒 自己想\n  - 1 分钟 跟 同桌 说\n  - 抽 2-3 个 学生 上台 分享\n• 关键 — Q6: 「如果 没有 ___ — 怎么样?」 戏剧化 想象\n• 老师 不 给标准答案 — 鼓励 多元 想法")

# 5-6. Transition slide
s=ns(); bg(s,HELP)
tb(s,1,1.10,8,0.80,"🎭",sz=100,a=PP_ALIGN.CENTER)
tb(s,1,2.15,8,0.70,"今天 — 我们 不只是 听 故事",sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,2.85,8,0.80,"我们 要 变成 — Community Helpers!",sz=32,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,1,3.75,8,0.40,"Today — we're not just listening to a story.",sz=14,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,4.15,8,0.40,"We're going to BECOME community helpers!",sz=14,b=True,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,4.85,8,0.40,"👩‍🏫  👩‍⚕️  👨‍🍳",sz=42,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎬 Transition · 30 秒:\n• 老师 戏剧化 地 念 这两行\n• 「今天 — 我们 都 要 演 一次 community helper!」\n• 升起 气氛 — 让 学生 兴奋 起来\n• 引出 — Part 3 角色扮演 开始!")

# ============================================================
# === SESSION 1 · PART 3: ROLE PLAY (30 min) ===
# ============================================================
s=phase_marker("🎭","Part 3 · Role Play","一日 小老师 / 小医生 / 小厨师",30,PH_PLAY,"全班 一起 — 几个 上台 演, 其他 配合 + 观察","Whole class — a few perform, the rest play along")
n+=1; pn(s,n)

# 6-1. Role Play rules
s=ns(); bg(s,CREAM); hb(s,"📋 Role Play 规则  Role Play Rules",PH_PLAY)
tb(s,0.4,0.85,9.2,0.40,"4 个 简单 规则 — 让 角色扮演 顺利 + 好玩!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"4 simple rules to make role play smooth + fun!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
rules=[
    ("1","🎬","进入 角色","Stay in character","当 你 是 老师 — 就像 老师 一样说话!",TEACH),
    ("2","🤝","互相 帮助","Help each other","队友 卡 住 — 你 帮 一下!",HELP),
    ("3","🙊","声音 适中","Voice volume","不 喊 不 闹 — 用 'inside voice'",PH_BOOK),
    ("4","🌟","试 一试","Try it","不会 没关系 — 试 就 是 英雄!",IDEA),
]
for i,(num,em,cn,en,detail,cl) in enumerate(rules):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.65+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.50),Inches(0.50))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+0.10,y+0.16,0.50,0.40,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.70,y+0.10,0.70,0.55,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+1.45,y+0.12,3.0,0.36,cn,sz=15,b=True,c=cl)
    tb(s,x+1.45,y+0.50,3.0,0.26,en,sz=9,c=GRAY)
    tb(s,x+0.15,y+0.92,4.30,0.50,detail,sz=11,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"📋 Rules · 2 分钟:\n• 老师 念 4 条 规则 — 用 戏剧 语气\n• 重点 — 第 4 条: 「不会 也 OK! 试 就 是 英雄!」\n• 让 K-5 都 安心 — 没有 完美 表演")

# 6-2. 小观察员任务
s=ns(); bg(s,CREAM); hb(s,"👀 小 观察员 任务  Observer Job",PH_PLAY)
tb(s,0.4,0.85,9.2,0.38,"不 上台 的 同学 — 你 也 有 任务!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.26,"Audience members — you have a job too!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
obs=[
    ("👂","用 耳朵 听","Listen with ears","他们 说 什么?",NAVY),
    ("👀","用 眼睛 看","Watch with eyes","他们 怎么 做?",HELP),
    ("💭","用 脑子 想","Think with brain","换 你 — 你 怎么 做?",PURPLE),
    ("👏","用 手 鼓掌","Hands clap","小帮手 — 你 真棒!",GOLD),
]
for i,(em,cn,en,detail,cl) in enumerate(obs):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y:=1.65),Inches(2.22),Inches(2.90))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.15,2.12,0.85,em,sz=50,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.05,2.12,0.40,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.45,2.12,0.28,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+1.80,2.02,1.0,detail,sz=11,c=DARK,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.78,"我 看到 ___ 在 ___ 。 我 觉得 ___ 。","I saw ___ doing ___. I think ___.",accent=PH_PLAY)
n+=1; pn(s,n)
notes(s,"👀 Observer · 1 分钟:\n• 强调: 「不 上台 也 是 重要 工作!」\n• 4 个 任务 — 耳朵 / 眼睛 / 脑子 / 手\n• 老师 每个 scenario 后 — 抽 1-2 个 观察员 分享 「我 看到 ___」\n• 让 全班 都 投入 — 没 人 闲着")

# ============================================================
# 7. STATION A — 👩‍🏫 一日小老师
# ============================================================
s=div("👩‍🏫 一日 小老师","Be a Teacher for a Day  ·  10 min  ·  4 scenarios",TEACH,"🎓"); n+=1; pn(s,n)

# 7-1. Teacher video
s=ns(); bg(s,CREAM); hb(s,"🎬 看 视频 — 真 老师 是 怎么 工作 的?",TEACH)
time_pill(s,2)
tb(s,0.4,0.85,9.2,0.34,"看 一段 真实 课堂 — 留意 老师 怎么 做!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Watch a real classroom — notice what the teacher does!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
vsh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.55),Inches(9.0),Inches(2.65))
vsh.fill.solid(); vsh.fill.fore_color.rgb=DARK; vsh.line.color.rgb=TEACH; vsh.line.width=Pt(3)
tb(s,0.5,1.85,9.0,1.40,"▶",sz=130,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.35,9.0,0.30,"👩‍🏫 真 老师 课堂 片段",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.5,3.70,9.0,0.30,"youtube.com/watch?v=HIcNSKZ65oU",sz=11,b=True,c=IDEA,a=PP_ALIGN.CENTER)
# Observation questions
obs=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.30),Inches(9.2),Inches(1.10))
obs.fill.solid(); obs.fill.fore_color.rgb=WARM; obs.line.color.rgb=TEACH; obs.line.width=Pt(2)
tb(s,0.55,4.38,9.0,0.28,"👀 看 的 时候 想 一想 (3 个 问题):",sz=12,b=True,c=TEACH)
tb(s,0.55,4.66,4.4,0.26,"1️⃣ 老师 在 做 什么?",sz=11,b=True,c=DARK)
tb(s,5.05,4.66,4.4,0.26,"2️⃣ 老师 帮助了 谁?",sz=11,b=True,c=DARK)
tb(s,0.55,4.98,9.0,0.26,"3️⃣ 她 是 怎么 解决 问题 的?",sz=11,b=True,c=DARK)
tb(s,0.55,5.26,9.0,0.22,"What is she doing? · Whom does she help? · How does she solve problems?",sz=8,c=GRAY)
n+=1; pn(s,n)
notes(s,"🎬 视频 · 2 分钟:\n• 视频: youtube.com/watch?v=HIcNSKZ65oU\n• 老师 先 念 3 个 观察问题, 再 播放\n• 看 完 — 不 讨论 — 直接 进入 4 个 scenario 演练!\n• 提示: 留意 老师 的 语气 / 手势 / 表情")

# 7-2 to 7-5. 4 Teacher Scenarios
teacher_scenarios=[
    ("🗣️","学生 一直 讲话","Students won't stop talking",
     "• 用 安静 信号 (举手 / 关灯)\n• 走 过去 — 不喊!\n• 说: 「请 安静 — 我们 来 听 ___」",
     ["请 安静。",
      "请 看 老师。",
      "我们 一起 听。"]),
    ("😢","有人 哭 了","Someone is crying",
     "• 蹲 下来 — 跟 学生 同 高度\n• 温柔 地 问 — 别 急\n• 听 — 不 急着 给 答案",
     ["你 怎么 了?",
      "我 来 帮 你。",
      "没 关系, 我 在 这里。"]),
    ("🤔","有 人 不会 做 题","Someone can't do the work",
     "• 不 直接 给 答案\n• 问: 「你 已经 知道 什么?」\n• 一步 一步 教 — 给 提示",
     ["我们 一起 试试。",
      "你 已经 知道 ___ 对吗?",
      "再 看 一下 这里。"]),
    ("🏃","大家 不会 排队","No one is lining up",
     "• 用 数 数 信号: 「3, 2, 1 — 排队!」\n• 表扬 第一个 排好 的\n• 用 唱 的 / 拍 节奏",
     ["请 排队。",
      "看 — 谁 排好 了?",
      "我们 一起 走。"]),
]
for em,scn_cn,scn_en,pro_action,frames in teacher_scenarios:
    s=ns(); bg(s,CREAM); hb(s,f"🎓 小老师 情景 — {scn_cn}",TEACH)
    scenario_card(s,1.0,em,scn_cn,scn_en,pro_action,frames,TEACH)
    tb(s,0.4,4.50,9.2,0.30,f"⏱ 演 2 分钟 · 1-2 个 学生 上台 当 小老师, 其他 当 学生 演 情景",sz=11,b=True,c=PH_PLAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.82,9.2,0.24,"1-2 students play the teacher; the rest play students in the scene",sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    sentence_frame_bar(s,5.05,"我 是 老师 — 我 会 ___ 。","I'm a teacher — I will ___.",accent=TEACH)
    n+=1; pn(s,n)
    notes(s,f"🎓 小老师 · {scn_cn} · 2 分钟:\n• 选 1 个 学生 当 小老师 (轮流 — 让 高/低 年级 都 试)\n• 3-4 个 学生 演 学生 (装 哭 / 装 不会 / 装 讲话...)\n• 1 分钟 演 + 1 分钟 反思: 「Professional 怎么 做?」\n• 老师 引导 — 不 评判 — 鼓励 试错\n• 课堂管理: 站 圆圈 — 全班 一起 看 — 不分 小组")

# ============================================================
# 8. STATION B — 👩‍⚕️ 一日小医生
# ============================================================
s=div("👩‍⚕️ 一日 小医生","Be a Doctor for a Day  ·  10 min  ·  4 scenarios",DOC,"🩺"); n+=1; pn(s,n)

doctor_scenarios=[
    ("🤢","肚子 疼","Stomach ache",
     "• 让 病人 坐 下\n• 温柔 地 问: 「哪里 不 舒服?」\n• 检查 + 不 着急",
     ["哪里 不 舒服?",
      "什么 时候 开始 疼?",
      "我 来 检查 一下。"]),
    ("🤕","摔倒 哭 了","Fell down and crying",
     "• 蹲 下来 — 看 伤口\n• 安慰: 「不 怕 — 我 在 这里」\n• 处理 伤口 — 慢慢 来",
     ["不 用 担心。",
      "我 来 看 一下。",
      "你 很 勇敢!"]),
    ("😨","害怕 打 针","Scared of shots",
     "• 不 假装 — 承认 「会 一点点 痛」\n• 给 选择: 「左手 还是 右手?」\n• 表扬 后 给 贴纸 / 棒棒糖",
     ["不 用 怕。",
      "深 呼吸 — 一 二 三!",
      "你 真 勇敢!"]),
    ("🪑","很多 病人 排队","Many patients waiting",
     "• 保持 微笑 — 不 慌\n• 一个 一个 看 — 不 跳过\n• 让 助手 帮 安抚 排队 的 人",
     ["请 稍 等 一下。",
      "我 马上 来。",
      "下 一位 — 请。"]),
]
for em,scn_cn,scn_en,pro_action,frames in doctor_scenarios:
    s=ns(); bg(s,CREAM); hb(s,f"🩺 小医生 情景 — {scn_cn}",DOC)
    scenario_card(s,1.0,em,scn_cn,scn_en,pro_action,frames,DOC)
    tb(s,0.4,4.50,9.2,0.30,f"⏱ 演 2 分钟 · 1-2 个 学生 上台 当 小医生, 其他 当 病人 演 情景",sz=11,b=True,c=PH_PLAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.82,9.2,0.24,"1-2 students play the doctor; the rest play patients in the scene",sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    sentence_frame_bar(s,5.05,"我 是 医生 — 我 会 ___ 。","I'm a doctor — I will ___.",accent=DOC)
    n+=1; pn(s,n)
    notes(s,f"🩺 小医生 · {scn_cn} · 2 分钟:\n• 道具: 玩具 听诊器 / 创可贴 / 纸杯 (病人喝水)\n• 选 1 个 学生 演 医生 (轮流)\n• 1-2 个 学生 演 病人 — 戏剧化 装 病\n• 关键 — 教 「同理心」: 蹲 下来 / 温柔 / 不 急\n• 高年级 可以 加 follow-up: 「你 多 大 了? 吃过 饭 了 吗?」")

# ============================================================
# 9. STATION C — 👨‍🍳 一日小厨师
# ============================================================
s=div("👨‍🍳 一日 小厨师","Be a Chef for a Day  ·  10 min  ·  4 scenarios",CHEF,"🍳"); n+=1; pn(s,n)

chef_scenarios=[
    ("🤤","客人 饿 了","Customer is hungry",
     "• 微笑 + 立即 招呼\n• 介绍 菜单 / 推荐\n• 快 但 不 慌",
     ["请 坐 — 您 想 吃 什么?",
      "我 推荐 ___。",
      "请 稍 等 — 马上 来!"]),
    ("🥲","食物 掉 地上 了","Food fell on the floor",
     "• 不 慌 — 道歉\n• 重 做 一份 — 不 给 客人\n• 学习: 下次 拿 稳",
     ["对 不起!",
      "我 再 做 一份。",
      "请 稍 等。"]),
    ("🥵","客人 说 太 辣 了","Customer says too spicy",
     "• 不 争 — 接受 意见\n• 给 一杯 凉水 / 牛奶\n• 重做 不 辣 的",
     ["对 不起 — 我 来 改。",
      "您 喝 一点 水。",
      "谢谢 您 的 意见。"]),
    ("🍴","餐厅 太 忙 了","Restaurant is too busy",
     "• 深 呼吸 — 不 慌\n• 跟 团队 一起 — 分工\n• 客人 等 — 道歉 + 谢谢",
     ["请 稍 等 — 我们 在 忙。",
      "谢谢 您 的 耐心。",
      "我 来 帮 你!"]),
]
for em,scn_cn,scn_en,pro_action,frames in chef_scenarios:
    s=ns(); bg(s,CREAM); hb(s,f"🍳 小厨师 情景 — {scn_cn}",CHEF)
    scenario_card(s,1.0,em,scn_cn,scn_en,pro_action,frames,CHEF)
    tb(s,0.4,4.50,9.2,0.30,f"⏱ 演 2 分钟 · 1-2 个 学生 上台 当 小厨师, 其他 当 客人 演 情景",sz=11,b=True,c=PH_PLAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.82,9.2,0.24,"1-2 students play the chef; the rest play customers in the scene",sz=8,c=GRAY,a=PP_ALIGN.CENTER)
    sentence_frame_bar(s,5.05,"我 是 厨师 — 我 会 ___ 。","I'm a chef — I will ___.",accent=CHEF)
    n+=1; pn(s,n)
    notes(s,f"🍳 小厨师 · {scn_cn} · 2 分钟:\n• 道具: 玩具 厨师 帽 / 围裙 / 假 食物 / 菜单 (打印)\n• 选 1-2 个 学生 演 厨师 (轮流)\n• 1-2 个 学生 演 客人\n• 关键 — 教 「服务 + 谦虚」: 不 争 / 道歉 / 改进\n• 高年级: 加 「点餐 → 做菜 → 上菜」 完整 流程\n• 低年级: 只 演 一个 简单 动作")

# ============================================================
# 10. SESSION 1 ENDING
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🌟 Session 1 总结  Community Helpers…",HELP)
tb(s,0.4,0.85,9.2,0.40,"今天 我们 体验 了 4 位 community helpers!",sz=18,b=True,c=HELP,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Today we experienced 4 community helpers!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3 takeaways
takes=[
    ("🤝","Help people","帮助 别人"),
    ("🛠️","Solve problems","解决 问题"),
    ("🏘","Make our community better","让 社区 更好"),
]
for i,(em,en,cn) in enumerate(takes):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(2.95),Inches(1.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=HELP; sh.line.color.rgb=IDEA; sh.line.width=Pt(3)
    tb(s,x+0.05,1.85,2.85,0.75,em,sz=50,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,2.85,0.40,cn,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,2.85,0.32,en,sz=11,c=IDEA,a=PP_ALIGN.CENTER)
# Reflection
refl=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.75),Inches(9.2),Inches(1.55))
refl.fill.solid(); refl.fill.fore_color.rgb=WARM; refl.line.color.rgb=HELP; refl.line.width=Pt(2)
tb(s,0.55,3.82,9.0,0.30,"💭 Reflection — 想 一想:",sz=14,b=True,c=HELP)
tb(s,0.55,4.15,9.0,0.30,"1️⃣ 哪个 工作 最 难?  Which job was hardest?",sz=12,b=True,c=DARK)
tb(s,0.55,4.50,9.0,0.30,"2️⃣ 哪个 工作 最 重要?  Which job is most important?",sz=12,b=True,c=DARK)
tb(s,0.55,4.85,9.0,0.30,"3️⃣ 你 今天 帮助 别人 了 吗?  Did YOU help someone today?",sz=12,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"🌟 Session 1 Close · 4-5 分钟:\n• 1 分钟: 全班 齐 喊: 「Community helpers help people, solve problems, make our community better!」\n• 3 分钟: Reflection — 选 1 题 全班 讨论\n  - 建议 Q3: 「你 今天 帮 了 谁?」 — 引出 第二天 课 (你 自己 也是 帮手!)\n• 1 分钟: Tease — 「下节课 — 复习 + 学新字 + 做 社区地图 + 颁奖!」\n• 课堂管理: 让 学生 站 起来 伸 个 懒腰 — 准备 第二节 课")

# ============================================================
# 11. SESSION 2 DIVIDER
# ============================================================
s=div("Session 2  ·  50 min","📖 复习 + 语言目标 + Projects",PH_LANG,"📚"); n+=1; pn(s,n)

# ============================================================
# === SESSION 2 · PART 1: REVIEW (15-20 min) ===
# ============================================================
s=phase_marker("🔄","Part 1 · Review","Session 1 回顾",18,PH_LANG,"配对 + 谁在哪里 + 情景回顾","Match + Where + Scenarios")
n+=1; pn(s,n)

# 12-1. Review game — 职业配对
s=ns(); bg(s,CREAM); hb(s,"🃏 复习 游戏 · 职业 配对",PH_LANG)
time_pill(s,5)
tb(s,0.4,0.85,9.2,0.36,"老师 说 一个 动作 — 哪一队 先 喊出 职业?",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Teacher names an action — first team to shout the job wins!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 8 action→job cards
actions=[
    ("✍️","写 黑板","→ 老师",TEACH),
    ("🩺","听 心跳","→ 医生",DOC),
    ("🍳","做 菜","→ 厨师",CHEF),
    ("🔥","救 火","→ 消防员",FIRE),
    ("📚","教 知识","→ 老师",TEACH),
    ("💊","开 药","→ 医生",DOC),
    ("🍰","烤 蛋糕","→ 厨师",CHEF),
    ("🚒","开 消防车","→ 消防员",FIRE),
]
for i,(em,action,answer,cl) in enumerate(actions):
    col=i%4; row=i//4
    x=0.4+col*2.32; y=1.55+row*1.75
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.22),Inches(1.62))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,x+0.05,y+0.08,2.12,0.70,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.82,2.12,0.36,action,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.20,2.12,0.32,answer,sz=12,b=True,c=cl,a=PP_ALIGN.CENTER)
tb(s,0.4,5.10,9.2,0.30,"🏆 第 一队 答对 +2 分! 4 队 抢答!",sz=12,b=True,c=PH_LANG,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🃏 配对 · 5 分钟 — 抢答:\n• 老师 盖住 答案 — 念 动作 — 各队 抢答\n• 答对 +2 分 (记 在 黑板 上)\n• 答错 不 扣分 — 下一组 接答\n• 课堂 管理: 用 举手 / 拍 桌子 抢答 — 不要 喊\n• 8 题 全部 做 完")

# 12-2. 谁在哪里工作?
s=ns(); bg(s,CREAM); hb(s,"📍 谁 在 哪里 工作?  Who Works Where?",PH_LANG)
time_pill(s,5)
tb(s,0.4,0.85,9.2,0.36,"用 一句话 说: 「___ 在 ___ 工作。」",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Use the sentence: '___ works at ___ .'",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# 4 person ↔ place pairs
pairs=[
    ("👩‍🏫","老师","→","🏫","学校",TEACH),
    ("👩‍⚕️","医生","→","🏥","医院",DOC),
    ("👨‍🍳","厨师","→","🍴","餐厅",CHEF),
    ("🧑‍🚒","消防员","→","🚒","消防局",FIRE),
]
for i,(p_em,p_cn,arrow,pl_em,pl_cn,cl) in enumerate(pairs):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.20,1.30,1.10,p_em,sz=50,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+1.10,1.30,0.32,p_cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+1.55,y+0.45,0.50,0.55,arrow,sz=24,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+2.15,y+0.20,2.30,1.10,pl_em,sz=50,a=PP_ALIGN.CENTER)
    tb(s,x+2.15,y+1.10,2.30,0.32,pl_cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.90,"___ 在 ___ 工作。","___ works at ___.",accent=PH_LANG)
n+=1; pn(s,n)
notes(s,"📍 谁在哪里 · 5 分钟:\n• 全班 跟读 4 对 — 慢/快/唱 3 遍\n• 抽 4 个 学生 上台 — 每人 说 一句\n• 鼓励 — 「错了 没关系 — 大声 说!」\n• Bonus: 加 第 5 对 — 「邮递员 → 邮局」 / 「警察 → 警察局」")

# 12-3. 情景回顾
s=ns(); bg(s,CREAM); hb(s,"🎬 情景 回顾  Scenario Replay",PH_LANG)
time_pill(s,7)
tb(s,0.4,0.85,9.2,0.36,"再 演 一次 — 你 最 喜欢 的 情景! ",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Replay your favorite scenario from yesterday!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Process steps
steps=[
    ("1","👥","2 人 一组","Pair up",PH_PLAY),
    ("2","🎯","选 1 个 情景","Pick a scenario",NAVY),
    ("3","🎬","演 1 分钟","Act 1 min",HELP),
    ("4","👏","换 队员 看 + 鼓掌","Switch + clap",GOLD),
]
for i,(num,em,cn,en,cl) in enumerate(steps):
    x=0.4+i*2.32
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.22),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.78),Inches(1.78),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+0.78,1.86,0.55,0.4,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.45,2.10,0.85,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.35,2.10,0.40,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.75,2.10,0.28,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
    if i<3:
        tb(s,x+2.10,2.85,0.30,0.4,"→",sz=22,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.85,"我 最 喜欢 的 是 ___ 情景, 因为 ___ 。","My favorite was ___ scenario, because ___.",accent=PH_PLAY)
n+=1; pn(s,n)
notes(s,"🎬 情景回顾 · 5-7 分钟:\n• 2 人 一组 (混龄 — 高/低 配)\n• 1 分钟 选 + 准备\n• 1 分钟 A 演 B 看\n• 1 分钟 换\n• 鼓励 高年级 帮 低年级\n• 老师 走动 — 听 + 表扬")

# ============================================================
# === SESSION 2 · PART 2: 我会认 / 我会写 (15 min) ===
# ============================================================
s=phase_marker("📝","Part 2 · 语言目标","我会认 + 我会写",15,PH_LANG,"5 字 认 + 2 字 写","5 read + 2 write")
n+=1; pn(s,n)

# 13. 我会认 — 5 words (compact gallery + per-word slides)
# Gallery first
s=ns(); bg(s,CREAM); hb(s,"📖 我 会 认 · 5 个 词  I Can Read · 5 Words",PH_LANG)
tb(s,0.4,0.85,9.2,0.36,"看 图 + 跟 老师 读!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
read_words=[
    ("📚","老师","lǎo shī",TEACH),
    ("🏫","学校","xué xiào",NAVY),
    ("🏥","医院","yī yuàn",DOC),
    ("🧑‍🚒","消防员","xiāo fáng yuán",FIRE),
    ("👨‍🍳","厨师","chú shī",CHEF),
]
for i,(em,cn,py,cl) in enumerate(read_words):
    col=i if i<3 else i-3+0
    x=0.4+(i%3)*3.10 if i<3 else 0.4+((i-3)+0.5)*3.10
    if i<3:
        x=0.4+i*3.10; y=1.35
    else:
        x=2.0+(i-3)*3.10; y=3.40
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.10,2.85,0.85,em,sz=54,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.00,2.85,0.46,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.50,2.85,0.28,py,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"📖 我会认 · gallery + 跟读 · 3-4 分钟:\n• 老师 指 一个 — 全班 齐 读 3 遍 (慢/快/唱)\n• 5 个 词 都 念 一遍\n• 检查: 「这 是 谁/什么?」 (指 emoji)\n• 引出 5 张 独立 slide — 例句")

# Per-word slides (compressed: just headline + sentence)
read_sentences=[
    ("📚","老师","lǎo shī","Teacher","老师 在 学校 工作。","Teachers work at school.",TEACH),
    ("🏫","学校","xué xiào","School","我 每天 去 学校。","I go to school every day.",NAVY),
    ("🏥","医院","yī yuàn","Hospital","医生 在 医院 工作。","Doctors work at the hospital.",DOC),
    ("🧑‍🚒","消防员","xiāo fáng yuán","Firefighter","消防员 救 火 — 真 勇敢!","Firefighters fight fires — brave!",FIRE),
    ("👨‍🍳","厨师","chú shī","Chef","厨师 会 做 饭。","Chefs can cook.",CHEF),
]
for em,cn,py,en,sent,sent_en,cl in read_sentences:
    s=ns(); bg(s,CREAM); hb(s,f"👀 我 会 认 · {cn}  I Can Read",cl)
    # Big character card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.5),Inches(2.6))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.fill.background()
    tb(s,0.5,1.10,4.3,1.4,cn,sz=56 if len(cn)==3 else 70,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.50,4.3,0.4,f"{py}  {en}",sz=16,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.95,4.3,0.4,"👉 跟 我 读!",sz=14,c=cl,a=PP_ALIGN.CENTER)
    # Big emoji card right
    ib=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.3),Inches(0.95),Inches(4.4),Inches(2.6))
    ib.fill.solid(); ib.fill.fore_color.rgb=WHITE; ib.line.color.rgb=cl; ib.line.width=Pt(3)
    tb(s,5.3,1.50,4.4,1.5,em,sz=120,a=PP_ALIGN.CENTER)
    # Sentence
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.3))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=cl; sh2.line.width=Pt(2)
    tb(s,0.6,3.95,1.5,0.4,"例 句",sz=14,b=True,c=cl)
    tb(s,0.6,4.30,8.8,0.5,sent,sz=22,b=True,c=DARK)
    tb(s,0.6,4.80,8.8,0.3,sent_en,sz=11,c=GRAY)
    n+=1; pn(s,n)
    notes(s,f"⏱ 1-2 分钟 — {cn}:\n• 指 字 + 全班 跟读 3 遍\n• 例句 念 1 遍 — 抽 1 个 学生 用 「{cn}」 造 新 句\n• 高年级 — 用 「{cn}」 写 一句")

# 14. 我会写 · 老师
s=ns(); bg(s,CREAM); hb(s,"✏️ 我 会 写 · 老师",TEACH)
tb(s,0.4,0.85,9.2,0.40,"一起 来 写 「老师」!",sz=22,b=True,c=TEACH,a=PP_ALIGN.CENTER)
tianzi(s,0.55,1.55,2.20,"老",TEACH,pinyin="lǎo",char_sz=120)
tianzi(s,2.95,1.55,2.20,"师",TEACH,pinyin="shī",char_sz=120)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.55),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=TEACH; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.55),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=TEACH; head.line.fill.background()
tb(s,5.45,1.62,4.10,0.4,"✏️ 怎么 写",sz=13,b=True,c=WHITE)
tb(s,5.45,2.20,4.10,0.40,"1️⃣ 「老」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.55,4.10,0.30,"  上「耂」+ 下「匕」",sz=10,c=GRAY)
tb(s,5.45,2.95,4.10,0.40,"2️⃣ 「师」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,3.30,4.10,0.30,"  左「丿+一」, 右像旗子!",sz=10,c=GRAY)
tb(s,5.45,3.75,4.10,0.40,"📝 田字格 写 3 遍",sz=12,b=True,c=TEACH)
sentence_frame_bar(s,4.55,"我 会 写 「老师」! 我 的 老师 叫 ___ 。","I can write 老师! My teacher's name is ___.",accent=TEACH)
n+=1; pn(s,n)
notes(s,"✏️ 老师 · 4-5 分钟:\n• 演示 笔顺 — 空中 写\n• 田字格 练 3 遍\n• 记忆: 「老」 上面 像 帽子 / 「师」 右边 像 旗子\n• 抽 学生 说: 「我的 老师 叫 ___」")

# 14b. 我会写 · 学校
s=ns(); bg(s,CREAM); hb(s,"✏️ 我 会 写 · 学校",NAVY)
tb(s,0.4,0.85,9.2,0.40,"一起 来 写 「学校」!",sz=22,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tianzi(s,0.55,1.55,2.20,"学",NAVY,pinyin="xué",char_sz=120)
tianzi(s,2.95,1.55,2.20,"校",NAVY,pinyin="xiào",char_sz=120)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.55),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=NAVY; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.55),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=NAVY; head.line.fill.background()
tb(s,5.45,1.62,4.10,0.4,"✏️ 怎么 写",sz=13,b=True,c=WHITE)
tb(s,5.45,2.20,4.10,0.40,"1️⃣ 「学」 — 8 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.55,4.10,0.30,"  上「⺍+冖」, 下「子」",sz=10,c=GRAY)
tb(s,5.45,2.95,4.10,0.40,"2️⃣ 「校」 — 10 笔",sz=14,b=True,c=DARK)
tb(s,5.45,3.30,4.10,0.30,"  左「木」 + 右「交」",sz=10,c=GRAY)
tb(s,5.45,3.75,4.10,0.40,"📝 田字格 写 3 遍",sz=12,b=True,c=NAVY)
sentence_frame_bar(s,4.55,"我 会 写 「学校」! 我 的 学校 叫 ___ 。","I can write 学校! My school is called ___.",accent=NAVY)
n+=1; pn(s,n)
notes(s,"✏️ 学校 · 4-5 分钟:\n• 演示 笔顺\n• 记忆: 「学」 上面 房盖, 下面 小孩 / 「校」 左边 木头\n• 田字格 练 3 遍\n• 抽 学生 说 自己 学校 的 名字")

# 14c. 句型练习
s=ns(); bg(s,CREAM); hb(s,"💬 句型 练习  Sentence Practice",PH_LANG)
time_pill(s,3)
tb(s,0.4,0.85,9.2,0.36,"用 3 个 句型 — 全班 跟读 + 自己 造句!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
sentences=[
    ("1","👩‍🏫","老师 在 学校 工作。","Teachers work at school.",TEACH),
    ("2","👩‍⚕️","医生 在 医院 工作。","Doctors work at the hospital.",DOC),
    ("3","👨‍🍳","厨师 会 做 饭。","Chefs can cook.",CHEF),
]
for i,(num,em,cn,en,cl) in enumerate(sentences):
    y=1.45+i*1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.0))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.55),Inches(y+0.20),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,0.55,y+0.28,0.55,0.40,num,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.10,1.0,0.80,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,2.50,y+0.18,7.0,0.40,cn,sz=20,b=True,c=DARK)
    tb(s,2.50,y+0.60,7.0,0.30,en,sz=11,c=GRAY)
sentence_frame_bar(s,4.95,"___ 在 ___ 工作。 / ___ 会 ___ 。","___ works at ___. / ___ can ___.",accent=PH_LANG)
n+=1; pn(s,n)
notes(s,"💬 句型 · 3 分钟:\n• 全班 跟读 3 遍\n• 抽 2-3 个 学生 — 用 句型 造 一个 新句\n• 鼓励 高年级 自己 编 — 低年级 跟着 念")

# ============================================================
# === SESSION 2 · PART 3: PROJECTS (15-20 min) ===
# ============================================================
s=phase_marker("🌟","Part 3 · Projects","Community Map + Thank You Awards",18,PH_PROJ,"P1 小组 地图 · P2 个人 奖状","P1 group map · P2 individual award")
n+=1; pn(s,n)

# ============================================================
# 15. PROJECT 1 — COMMUNITY MAP
# ============================================================
# 15-1. Project 1 intro
s=ns(); bg(s,CREAM); hb(s,"🏘 Project 1 · 社区 地图  Community Map",PH_PROJ)
time_pill(s,10)
tb(s,0.4,0.85,9.2,0.42,"小组 一起 — 画 一张 你们 的 「社区 地图」!",sz=18,b=True,c=PH_PROJ,a=PP_ALIGN.CENTER)
tb(s,0.4,1.32,9.2,0.28,"In teams — draw your community map together!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 5 places they should include
places=[
    ("🏫","学校"),("🏥","医院"),("🍴","餐厅"),("🚒","消防局"),("🌳","公园"),
]
for i,(em,cn) in enumerate(places):
    x=0.4+i*1.90
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(1.80),Inches(1.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=PH_PROJ; sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.85,1.70,0.75,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.65,1.70,0.40,cn,sz=15,b=True,c=PH_PROJ,a=PP_ALIGN.CENTER)
# 4 helpers to add inside
banner=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.45),Inches(9.2),Inches(1.05))
banner.fill.solid(); banner.fill.fore_color.rgb=HELP; banner.line.color.rgb=IDEA; banner.line.width=Pt(2.5)
tb(s,0.55,3.52,9.0,0.30,"➕ 并且 在 每个 地方 加上 一位 小帮手:",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.55,3.85,9.0,0.30,"And — add a helper to each place:",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.55,4.18,9.0,0.30,"👩‍🏫 老师  ·  👩‍⚕️ 医生  ·  👨‍🍳 厨师  ·  🧑‍🚒 消防员",sz=14,b=True,c=IDEA,a=PP_ALIGN.CENTER)
tb(s,0.4,4.70,9.2,0.30,"⏱ 10 分钟 — 画 + 标 名字 + 加 小帮手",sz=12,b=True,c=PH_PROJ,a=PP_ALIGN.CENTER)
tb(s,0.4,5.05,9.2,0.26,"10 min — draw + label + add helpers",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🏘 Project 1 Intro · 2 分钟:\n• 4-5 人 一组 (混龄 — 老师 分配)\n• 每组 1 张 大白纸 + marker / 蜡笔\n• 任务 — 画 一个 社区, 包括 5 个 地方 + 4 位 小帮手\n• 强调 — 不用 完美 — 大胆 画!")

# 15-2. Example map
s=ns(); bg(s,CREAM); hb(s,"🗺️ 例子 — 像 这样 画!  Example Map",PH_PROJ)
tb(s,0.4,0.85,9.2,0.34,"5 个 地方 + 4 位 小帮手 — 都 在 一张 纸 上!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"5 places + 4 helpers — all on one paper!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Visual example — 5 buildings with arrows + helpers around
example_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.55),Inches(9.0),Inches(3.30))
example_box.fill.solid(); example_box.fill.fore_color.rgb=WARM; example_box.line.color.rgb=PH_PROJ; example_box.line.width=Pt(3)
# Sun
tb(s,0.6,1.65,0.8,0.7,"☀️",sz=42,a=PP_ALIGN.CENTER)
# Place blocks
place_layout=[
    (1.40,2.20,"🏫","学校","👩‍🏫"),
    (3.30,2.20,"🏥","医院","👩‍⚕️"),
    (5.20,2.20,"🍴","餐厅","👨‍🍳"),
    (7.10,2.20,"🚒","消防局","🧑‍🚒"),
    (4.30,3.55,"🌳","公园","👶"),
]
for x,y,p_em,pl_cn,h_em in place_layout:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(1.50),Inches(1.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=PH_PROJ; sh.line.width=Pt(2)
    tb(s,x+0.05,y+0.05,1.40,0.50,p_em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.55,1.40,0.28,pl_cn,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.80,1.40,0.30,h_em,sz=20,a=PP_ALIGN.CENTER)
# Roads (dashed look) — simple line shapes
for y_road in [3.40]:
    rd=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.5),Inches(y_road),Inches(7.10),Inches(0.05))
    rd.fill.solid(); rd.fill.fore_color.rgb=GRAY; rd.line.fill.background()
tb(s,0.4,5.05,9.2,0.30,"💡 加 你 自己的 想法 — 比如 警察局、邮局、超市……",sz=12,b=True,c=PH_PROJ,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🗺️ 例子 · 1 分钟:\n• 老师 指 例子 — 不 念 太 久 — 让 学生 自己 想象\n• 鼓励 — 加 自己 喜欢 的 地方 (超市? 公园? 海边?)")

# 15-3. Materials + teamwork
s=ns(); bg(s,CREAM); hb(s,"🛠️ 材料 + 团队 规则  Materials + Team Rules",PH_PROJ)
# Materials box (left)
mat=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.55),Inches(4.20))
mat.fill.solid(); mat.fill.fore_color.rgb=WHITE; mat.line.color.rgb=PH_PROJ; mat.line.width=Pt(3)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.55),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=PH_PROJ; head.line.fill.background()
tb(s,0.55,1.03,4.30,0.4,"🛠️ 材料  Materials",sz=16,b=True,c=WHITE)
mats=[("📄","大 白纸 (每组 1 张)","Large paper (1/group)"),
      ("🖍️","蜡笔 / Marker","Crayons / Markers"),
      ("✂️","剪刀 + 胶水 (可选)","Scissors + glue (opt)"),
      ("📌","便利贴 (写 名字)","Sticky notes")]
for i,(em,cn,en) in enumerate(mats):
    y=1.65+i*0.78
    tb(s,0.60,y,0.5,0.45,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.05,3.65,0.36,cn,sz=12,b=True,c=DARK)
    tb(s,1.20,y+0.40,3.65,0.26,en,sz=9,c=GRAY)
# Teamwork rules (right)
tm=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.95),Inches(4.55),Inches(4.20))
tm.fill.solid(); tm.fill.fore_color.rgb=WHITE; tm.line.color.rgb=HELP; tm.line.width=Pt(3)
head2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.95),Inches(4.55),Inches(0.50))
head2.fill.solid(); head2.fill.fore_color.rgb=HELP; head2.line.fill.background()
tb(s,5.20,1.03,4.30,0.4,"👥 团队 规则  Team Rules",sz=16,b=True,c=WHITE)
rules=[("1️⃣","每个 人 都 画!","Everyone draws"),
       ("2️⃣","互相 帮助 — 不 抢","Help, don't grab"),
       ("3️⃣","声音 小 + 听 队友","Soft voice + listen"),
       ("4️⃣","时间 到 — 一起 分享","Share when time's up")]
for i,(em,cn,en) in enumerate(rules):
    y=1.65+i*0.78
    tb(s,5.25,y,0.5,0.45,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,5.85,y+0.05,3.65,0.36,cn,sz=12,b=True,c=DARK)
    tb(s,5.85,y+0.40,3.65,0.26,en,sz=9,c=GRAY)
n+=1; pn(s,n)
notes(s,"🛠️ 材料 + 规则 · 1 分钟:\n• 老师 提前 准备 好 材料 — 每组 一套\n• 强调 团队 规则 — 尤其 第 1 条 (每人 都画) + 第 2 条 (不抢)\n• 老师 课堂 管理: 站 在 教室 中间 — 4 队 都能看到\n• 老师 走动 — 拍 1-2 张 照片 — 记录 学生 作品")

# 15-4. Sharing
s=ns(); bg(s,CREAM); hb(s,"🎤 分享 时间  Map Sharing",PH_PROJ)
time_pill(s,3)
tb(s,0.4,0.85,9.2,0.36,"每 队 派 1 个 代表 — 上台 介绍 1 分钟!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.26,"Each team — send 1 rep to share for 1 min!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3 discussion questions
disc=[
    ("🔍","谁 在 帮助 大家?","Who's helping?",HELP),
    ("⭐","哪个 地方 最 重要?","Most important place?",NAVY),
    ("😱","如果 没有 医院 — 会 怎样?","If no hospital — what happens?",DOC),
]
for i,(em,cn,en,cl) in enumerate(disc):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.95),Inches(2.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+1.20),Inches(1.78),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+1.20,1.86,0.55,0.4,str(i+1),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.50,2.85,0.85,em,sz=46,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,2.85,0.40,cn,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.85,2.85,0.28,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.65,"我们 的 社区 有 ___, ___ 在 帮助 大家。","Our community has ___, where ___ helps people.",accent=PH_PROJ)
n+=1; pn(s,n)
notes(s,"🎤 分享 · 3 分钟:\n• 每队 60 秒 — 用 句型 介绍\n• 全班 鼓掌 — 每队 都 +5 分\n• 选 1 个 讨论 问题 — 抽 2 个 学生 回答\n• 关键 — Q3: 「如果 没有 医院 / 学校 — 会怎样?」 强化 「 帮手 重要 」")

# ============================================================
# 16. PROJECT 2 — THANK YOU AWARDS
# ============================================================
# 16-1. Project 2 intro
s=ns(); bg(s,CREAM); hb(s,"🏅 Project 2 · Thank You Awards  感谢 奖状",PH_PROJ)
time_pill(s,7)
tb(s,0.4,0.85,9.2,0.42,"给 一位 community helper 颁 奖 — 你 来 设计 奖状!",sz=17,b=True,c=PH_PROJ,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Design a thank-you award for a community helper!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 4 award examples
awards=[
    ("🏆","最佳 老师 奖","Best Teacher Award",TEACH),
    ("🥇","超级 医生 奖","Super Doctor Award",DOC),
    ("❤️","爱心 厨师 奖","Caring Chef Award",CHEF),
    ("🦸","勇敢 消防员 奖","Brave Firefighter Award",FIRE),
]
for i,(em,cn,en,cl) in enumerate(awards):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.65+row*1.60
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    tb(s,x+0.10,y+0.15,1.20,1.10,em,sz=48,a=PP_ALIGN.CENTER)
    tb(s,x+1.40,y+0.18,3.0,0.45,cn,sz=18,b=True,c=cl)
    tb(s,x+1.40,y+0.68,3.0,0.30,en,sz=10,c=GRAY)
    tb(s,x+1.40,y+1.00,3.0,0.30,"⭐ 颁给 一位 ___!",sz=11,b=True,c=DARK)
tb(s,0.4,4.95,9.2,0.30,"或者 — 你 自己 设计 新 奖状! Or — design your OWN award!",sz=12,b=True,c=PH_PROJ,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🏅 Project 2 Intro · 1 分钟:\n• 个人 项目 — 每人 1 张 奖状 纸\n• 4 个 example — 启发 — 不一定 要 follow\n• 高年级 可以 设计 新 奖项\n• 颁给 谁? 真人 (你妈 / 你老师) 都 可以!")

# 16-2. Award template + sentence frame
s=ns(); bg(s,CREAM); hb(s,"📝 写 一张 奖状  Fill in Your Award",PH_PROJ)
tb(s,0.4,0.85,9.2,0.34,"填 4 个 空 + 画 一个 图 + 签 你 的 名字!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.18,9.2,0.24,"Fill in 4 blanks + draw + sign your name!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Award template
award_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.55),Inches(9.0),Inches(3.45))
award_box.fill.solid(); award_box.fill.fore_color.rgb=WARM; award_box.line.color.rgb=GOLD; award_box.line.width=Pt(5)
# Title bar
title_bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.7),Inches(1.70),Inches(8.6),Inches(0.55))
title_bar.fill.solid(); title_bar.fill.fore_color.rgb=GOLD; title_bar.line.fill.background()
tb(s,0.7,1.78,8.6,0.4,"🏅 ___________________ 奖",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Body
tb(s,0.85,2.40,8.30,0.40,"颁 给:  ____________________________",sz=18,b=True,c=DARK)
tb(s,0.85,2.85,8.30,0.40,"因为 你:  ___________________________",sz=18,b=True,c=DARK)
tb(s,0.85,3.30,8.30,0.40,"我 想 说:  谢谢 你 ___________________",sz=18,b=True,c=DARK)
tb(s,0.85,3.75,8.30,0.40,"我 的 名字:  _______________________",sz=18,b=True,c=DARK)
tb(s,0.85,4.30,8.30,0.30,"🎨 在 这里 画 一个 小 图",sz=12,b=True,c=GOLD)
sentence_frame_bar(s,5.10,"谢谢 ___ , 因为 你 帮助 别人。","Thank you ___, because you help others.",accent=GOLD)
n+=1; pn(s,n)
notes(s,"📝 奖状 · 5-6 分钟:\n• 老师 发 奖状 纸 (打印 / 手画 模板)\n• 5 分钟 写 + 画\n• 老师 走动 — 帮 拼字 / 鼓励\n• 高年级 — 句子 要 完整\n• 低年级 — 几个 字 + 多 画图\n• 写 完 — 不收 — 留 着 颁奖典礼 用!")

# 16-3. Award Ceremony
s=ns(); bg(s,GOLD)
tb(s,1,0.50,8,0.85,"🎉",sz=100,a=PP_ALIGN.CENTER)
tb(s,1,1.55,8,0.80,"Community Helper",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,2.30,8,0.80,"Award Ceremony!",sz=36,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.55,"🏅 颁奖 典礼 🏅",sz=28,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Process
proc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(3.95),Inches(8.0),Inches(1.20))
proc.fill.solid(); proc.fill.fore_color.rgb=WHITE; proc.line.color.rgb=DARK; proc.line.width=Pt(3)
tb(s,1.0,4.02,8.0,0.30,"⏱ 5 分钟 — 流程:",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1.0,4.32,8.0,0.30,"1️⃣ 上台 念 奖状  2️⃣ 全班 鼓掌  3️⃣ 老师 拍照",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1.0,4.65,8.0,0.30,"4️⃣ 把 奖状 真的 送 给 ta!",sz=14,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎉 颁奖典礼 · 5 分钟:\n• 老师 当 主持人 — 戏剧化\n• 每个 学生 30 秒 上台 — 念 自己 奖状\n• 全班 大力 鼓掌\n• 关键 — 「真的 送出去!」 — 学生 回家 / 课后 给 ta\n• 老师 拍 集体 合影 — 这 是 大家 的 时刻!\n• 课堂 管理: 让 学生 排队 — 一个 一个 上 — 别 抢")

# ============================================================
# 17. DAY 4 BADGE — CLOSING
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🎖️ Day 4 完成!  Day 4 Complete!",PH_CLOSE)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(0.95),Inches(3),Inches(3))
sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=HEART; sh.line.width=Pt(6)
tf=tb(s,3.6,1.20,2.8,2.7,"DAY 4",sz=18,b=True,c=HEART,a=PP_ALIGN.CENTER)
ap(tf,"🏘 ❤️",sz=46,a=PP_ALIGN.CENTER)
ap(tf,"社区小帮手",sz=16,b=True,c=HELP,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=12,b=True,c=OK,a=PP_ALIGN.CENTER)
ap(tf,"👩‍🏫 👩‍⚕️ 👨‍🍳 🧑‍🚒",sz=18,a=PP_ALIGN.CENTER)
tb(s,0.4,4.20,9.2,0.40,"🎉 你 今天 帮助 别人 了! 你 也 是 一位 community helper!",sz=15,b=True,c=HELP,a=PP_ALIGN.CENTER)
tb(s,0.4,4.65,9.2,0.30,"You helped others today! You are a community helper too!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tease=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.05),Inches(9.2),Inches(0.40))
tease.fill.solid(); tease.fill.fore_color.rgb=PURPLE; tease.line.fill.background()
tb(s,0.55,5.10,9.0,0.32,"🤖 明天 Day 5 — AI 与未来! 机器人 也是 community helper 吗?",sz=12,b=True,c=IDEA,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎖️ 收尾 · 2 分钟:\n• 全班 齐 喊: 「Community helpers help, solve, make better! 我 也是!」\n• 发 徽章 / 贴纸\n• 预告 Day 5 — AI 主题")

# ============================================================
# CLASSROOM MANAGEMENT TIPS (附录)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📋 老师 备忘  Classroom Management Tips",NAVY)
tb(s,0.4,0.85,9.2,0.32,"K-5 混龄 · 约 20 人 · 高 互动 课堂 的 关键",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tips=[
    ("👥","分组 策略","Grouping","4-5 人 一组 · 混龄 (高/低 配) · 老师 提前 分好",NAVY),
    ("🎭","Role play","Stations","站 圆圈 / 教室 中间 · 4 队 轮流 · 不 抢 不 闹",PH_PLAY),
    ("⏰","时间 控制","Timing","用 计时器 + 手势 · 提前 30 秒 提醒",GOLD),
    ("🤝","混龄 配合","Mixed-age","高年级 当 mentor · 低年级 当 学习者 · 互相 帮助",HELP),
    ("🙊","声音 管理","Voice","用 安静 信号 (举手 / 关灯) — 不 喊!",PURPLE),
    ("📸","记录 + 表扬","Capture","拍 照 / 视频 — 课后 给 家长 看 + 多 表扬",HEART),
]
for i,(em,cn,en,detail,cl) in enumerate(tips):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.30+row*1.30
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.15))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,x+0.10,y+0.15,0.80,0.85,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+1.05,y+0.10,3.40,0.36,cn,sz=14,b=True,c=cl)
    tb(s,x+1.05,y+0.46,3.40,0.26,en,sz=9,c=GRAY)
    tb(s,x+1.05,y+0.75,3.40,0.35,detail,sz=10,b=True,c=DARK)
n+=1; pn(s,n)
notes(s,"📋 老师 备忘 (不 在 课堂 投影 — 只是 老师 参考):\n• 课前 准备:\n  - 4 队 名单 (混龄)\n  - role play 道具 (听诊器 / 厨师帽 / 粉笔...)\n  - 视频 链接 + 投影 测试\n  - 奖状 纸 / 大白纸 / marker\n  - 计时器\n• 课中:\n  - 抓 节奏 — 不 卡 在 任何 一个 scenario\n  - 高年级 当 mentor — 不 全 包办\n  - 老师 走动 — 不 站 讲台\n• 课后:\n  - 拍 照 给 家长\n  - 留 学生 作品 (奖状 / 地图)\n  - 短 反思: 哪个 部分 学生 最 投入?")

# ============================================================
out=os.path.join(os.path.dirname(__file__),"day4_helpers.pptx")
prs.save(out); print(f"Saved {out}  ({n} slides)")
