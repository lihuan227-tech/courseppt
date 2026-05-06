#!/usr/bin/env python3
"""
我的职业梦想 — Day 2: 小小科学家 / 工程师 (Problem Solver Day)
3 sessions × 50 min, Workshop Model (5-phase frame per session):
  🔥 Hook (5) → 📚 Mini-Lesson (10) → 🎯 Active Practice (20) → 🌱 Apply (10) → 🎤 Share & Close (5)
Adapted for 大班 (15+) with 分组比赛 (4 teams) + 轮流上台.
4 problem solvers: 爱迪生 · 简·古道尔 · 莱特兄弟 · Garrett Morgan
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# === Palette ===
SCI    = RGBColor(0x00,0x69,0x6B)   # scientist teal
ENG    = RGBColor(0xC4,0x52,0x2A)   # engineer rust
YELLOW = RGBColor(0xF5,0xC2,0x42)   # curiosity gold
NAVY   = RGBColor(0x1E,0x3A,0x5F)
LAB    = RGBColor(0xE5,0x3E,0x3E)   # alarm/problem red
GREEN  = RGBColor(0x2E,0x7D,0x32)
PURPLE = RGBColor(0x7B,0x1F,0xA2)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
BROWN  = RGBColor(0x6B,0x44,0x23)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
OK     = RGBColor(0x38,0x8E,0x3C)

# Phase colors (consistent across sessions)
PH_HOOK   = LAB     # 🔥 high energy
PH_MINI   = SCI     # 📚 calm focus
PH_ACTIVE = ENG     # 🎯 highest energy
PH_APPLY  = GREEN   # 🌱 individual focus
PH_CLOSE  = YELLOW  # 🎤 celebration

# Per-person accent colors
EDISON  = RGBColor(0xF5,0xA6,0x23)  # bulb amber
GOODALL = RGBColor(0x4C,0x6E,0x4F)  # forest
WRIGHT  = RGBColor(0x1F,0x77,0xB4)  # sky blue
MORGAN  = RGBColor(0xD3,0x35,0x35)  # traffic red

# Team colors (持续使用 — 4 teams across all 3 sessions)
T_RED    = RGBColor(0xE5,0x3E,0x3E)
T_BLUE   = RGBColor(0x1F,0x77,0xB4)
T_GREEN  = RGBColor(0x2E,0x7D,0x32)
T_YELLOW = RGBColor(0xF5,0xA6,0x23)

# === Helpers ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def hb(s,txt,c=SCI,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,text):
    nf=s.notes_slide.notes_text_frame; lines=text.split("\n"); nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line
def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    return s
def sentence_frame_bar(s,t,frame_cn,frame_en,accent=YELLOW):
    if t > 4.95: t = 4.95
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid(); sf.fill.fore_color.rgb=WARM; sf.line.color.rgb=accent; sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.7,0.4,"💬 我来说",sz=14,b=True,c=accent)
    tb(s,2.0,t+0.07,7.6,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.0,t+0.32,7.6,0.3,frame_en,sz=10,c=GRAY)

# Phase marker — full slide intro for each phase (10-15 sec)
def phase_marker(emoji,phase_cn,phase_en,time_min,color,what_cn,what_en):
    s=ns(); bg(s,color)
    tb(s,1,0.85,8,0.7,emoji,sz=80,a=PP_ALIGN.CENTER)
    tb(s,1,1.85,8,0.6,phase_cn,sz=38,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.50,8,0.4,phase_en,sz=16,c=YELLOW,a=PP_ALIGN.CENTER)
    # Time pill
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.5),Inches(3.10),Inches(3.0),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=YELLOW; sh.line.fill.background()
    tb(s,3.5,3.18,3.0,0.4,f"⏱  {time_min} 分钟",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # What we'll do
    tb(s,1,4.00,8,0.5,what_cn,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,4.55,8,0.35,what_en,sz=12,c=YELLOW,a=PP_ALIGN.CENTER)
    return s

# Score badge — top-right corner widget showing 4 teams
def score_badge(s):
    teams=[("🔴",T_RED),("🔵",T_BLUE),("🟢",T_GREEN),("🟡",T_YELLOW)]
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.40),Inches(0.78),Inches(2.30),Inches(0.32))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=GRAY; sh.line.width=Pt(0.75)
    tb(s,7.45,0.81,0.45,0.28,"🏆",sz=12,a=PP_ALIGN.CENTER)
    for i,(em,cl) in enumerate(teams):
        tb(s,7.85+i*0.45,0.81,0.40,0.28,f"{em}__",sz=10,b=True,c=cl,a=PP_ALIGN.CENTER)

# Group task label — small badge top-left
def group_label(s,t=0.78):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.30),Inches(t),Inches(1.80),Inches(0.32))
    sh.fill.solid(); sh.fill.fore_color.rgb=PH_ACTIVE; sh.line.fill.background()
    tb(s,0.30,t+0.03,1.80,0.28,"👥 分组任务",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)

# TPR cue strip
def tpr_strip(s,t,cue_cn,cue_en):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=PH_HOOK; sh.line.fill.background()
    tb(s,0.5,t+0.05,9.0,0.3,f"🙌 TPR · {cue_cn}",sz=14,b=True,c=WHITE)
    tb(s,0.5,t+0.30,9.0,0.25,cue_en,sz=10,c=YELLOW)

# Vote strip — for "科学家 还是 工程师?" categorize moments
def vote_strip(s,t,answer):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=PURPLE; sh.line.fill.background()
    tb(s,0.5,t+0.05,5.0,0.3,"🗳️ 小组投票  Team Vote",sz=14,b=True,c=WHITE)
    tb(s,0.5,t+0.30,5.0,0.25,"科学家 (蓝牌) 还是 工程师 (红牌)?",sz=10,c=YELLOW)
    # Answer pill
    pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.80),Inches(t+0.10),Inches(3.7),Inches(0.36))
    pill.fill.solid(); pill.fill.fore_color.rgb=YELLOW; pill.line.fill.background()
    tb(s,5.80,t+0.13,3.7,0.30,f"答: {answer}  +1 分",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)

# Question pill (single-line question card)
def question_pill(s,x,y,w,h,em,q_cn,q_en,color):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=color; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.10,0.55,h-0.2,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.75,y+0.08,w-0.85,0.40,q_cn,sz=15,b=True,c=DARK)
    tb(s,x+0.75,y+0.45,w-0.85,0.30,q_en,sz=10,c=GRAY)

# VIDEO slide — dedicated video viewing slide before each story
def video_slide(emoji,cn,en,color,video_cn,video_en,watch_for_cn="",watch_for_en=""):
    """Dedicated video slide — large placeholder + viewing hints + what-to-watch-for prompt."""
    s=ns(); bg(s,CREAM); hb(s,f"{emoji} 先看视频  Watch Video · {cn}  ·  {en}",color)
    score_badge(s)
    # Large 16:9 video placeholder
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(0.95),Inches(9.0),Inches(3.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=DARK
    sh.line.color.rgb=color; sh.line.width=Pt(3)
    tb(s,0.5,1.50,9.0,1.5,"▶",sz=130,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.45,9.0,0.30,"老师在这里插入视频 / Teacher: insert video here",sz=12,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
    # Video search hints strip (0.70" tall to fit 2-line CN + EN)
    vh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.20),Inches(9.2),Inches(0.70))
    vh.fill.solid(); vh.fill.fore_color.rgb=PURPLE; vh.line.color.rgb=YELLOW; vh.line.width=Pt(1.5)
    tb(s,0.55,4.34,2.0,0.30,"🔍 视频建议",sz=12,b=True,c=YELLOW)
    tb(s,2.55,4.26,7.0,0.42,video_cn,sz=11,b=True,c=WHITE)
    tb(s,2.55,4.69,7.0,0.20,video_en,sz=8,c=WARM)
    # What to watch for (optional viewing prompt)
    if watch_for_cn:
        wf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.00),Inches(9.2),Inches(0.55))
        wf.fill.solid(); wf.fill.fore_color.rgb=WARM; wf.line.color.rgb=color; wf.line.width=Pt(1.5)
        tb(s,0.55,5.05,9.0,0.22,f"👀 看的时候想想…  {watch_for_cn}",sz=12,b=True,c=color)
        tb(s,0.55,5.27,9.0,0.18,watch_for_en,sz=9,c=GRAY)
    return s

# STORY slide — answer 3 questions (with expected answers as body text)
def story_slide(emoji,cn,en,years,country,color,panels,fun_cn,fun_en,video_cn="",video_en=""):
    """panels: list of 3 tuples (panel_emoji, q_cn, q_en, ans_cn, ans_en, pause_cn)
    video_cn/en: legacy params, retained for back-compat (now rendered on a separate video_slide).
    """
    s=ns(); bg(s,CREAM); hb(s,f"{emoji} 故事 · {cn}  ·  Story of {en}",color)
    score_badge(s)
    # Identity strip below header
    idstrip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(0.50))
    idstrip.fill.solid(); idstrip.fill.fore_color.rgb=color; idstrip.line.fill.background()
    tb(s,0.55,0.90,0.7,0.40,emoji,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.30,0.90,3.5,0.35,cn,sz=17,b=True,c=WHITE)
    tb(s,5.00,0.92,2.4,0.22,f"📅 {years}",sz=10,c=YELLOW)
    tb(s,5.00,1.13,2.4,0.20,en,sz=9,c=YELLOW)
    tb(s,7.45,0.92,2.0,0.22,f"📍 {country}",sz=10,c=YELLOW)
    # 3-panel comic — each panel = 1 question + expected answer + pause prompt
    panel_y=1.50
    panel_h=3.10
    for i,(p_em,p_cn,p_en,b_cn,b_en,pause_cn) in enumerate(panels):
        x=0.4+i*3.10
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(panel_y),Inches(2.95),Inches(panel_h))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=color; sh.line.width=Pt(2.5)
        # Number badge top-left
        nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(panel_y+0.07),Inches(0.42),Inches(0.42))
        nb.fill.solid(); nb.fill.fore_color.rgb=color; nb.line.fill.background()
        tb(s,x+0.10,panel_y+0.13,0.42,0.32,str(i+1),sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        # Big panel emoji
        tb(s,x+0.55,panel_y+0.07,2.4,0.50,p_em,sz=28,a=PP_ALIGN.CENTER)
        # Question heading
        tb(s,x+0.10,panel_y+0.62,2.85,0.32,p_cn,sz=14,b=True,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.10,panel_y+0.92,2.85,0.22,p_en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
        # Expected answer
        tb(s,x+0.15,panel_y+1.20,2.75,0.45,b_cn,sz=11,c=DARK)
        tb(s,x+0.15,panel_y+1.62,2.75,0.30,b_en,sz=8,c=GRAY)
        # Pause prompt at bottom of panel
        ps_y=panel_y+panel_h-0.43
        ps=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.10),Inches(ps_y),Inches(2.75),Inches(0.38))
        ps.fill.solid(); ps.fill.fore_color.rgb=YELLOW; ps.line.fill.background()
        tb(s,x+0.18,ps_y+0.04,2.65,0.30,f"👉 {pause_cn}",sz=10,b=True,c=DARK)
    # Fun fact strip at bottom
    ff_y=panel_y+panel_h+0.05
    ff=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(ff_y),Inches(9.2),Inches(0.45))
    ff.fill.solid(); ff.fill.fore_color.rgb=WARM; ff.line.color.rgb=color; ff.line.width=Pt(1.5)
    tb(s,0.55,ff_y+0.04,0.6,0.40,"⭐",sz=16,a=PP_ALIGN.CENTER)
    tb(s,1.10,ff_y+0.04,8.4,0.25,fun_cn,sz=11,b=True,c=color)
    tb(s,1.10,ff_y+0.27,8.4,0.18,fun_en,sz=8,c=GRAY)
    return s

# Q&A slide — 2 comprehension questions + vote + TPR + sentence frame
def qa_slide(emoji,cn,en,color,questions,vote_answer,tpr_cn,tpr_en,sentence_cn,sentence_en):
    """questions: list of 2 tuples (q_emoji, q_cn, q_en, hint_cn, hint_en)"""
    s=ns(); bg(s,CREAM); hb(s,f"{emoji} 互动问答 · {cn}  ·  Q&A",color)
    group_label(s)
    score_badge(s)
    # 2 question pills
    for i,(q_em,q_cn,q_en,hint_cn,hint_en) in enumerate(questions):
        y=1.20+i*0.95
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.85))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=color; sh.line.width=Pt(2)
        tb(s,0.55,y+0.18,0.6,0.5,q_em,sz=24,a=PP_ALIGN.CENTER)
        tb(s,1.30,y+0.10,5.3,0.35,q_cn,sz=14,b=True,c=DARK)
        tb(s,1.30,y+0.45,5.3,0.30,q_en,sz=10,c=GRAY)
        # Hint pill (right side)
        pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.70),Inches(y+0.15),Inches(2.80),Inches(0.55))
        pill.fill.solid(); pill.fill.fore_color.rgb=YELLOW; pill.line.fill.background()
        tb(s,6.75,y+0.18,2.70,0.28,f"💡 {hint_cn}",sz=11,b=True,c=DARK)
        tb(s,6.75,y+0.43,2.70,0.22,hint_en,sz=9,c=DARK)
    # Vote strip
    vote_strip(s,3.20,vote_answer)
    # TPR strip
    tpr_strip(s,3.85,tpr_cn,tpr_en)
    # Sentence frame
    sentence_frame_bar(s,4.50,sentence_cn,sentence_en,accent=color)
    return s

n=0

# ============================================================
# 1. COVER
# ============================================================
s=ns(); bg(s,CREAM)
tb(s,1,0.40,8,0.55,"我的职业梦想 · My Dream Career",sz=22,b=True,c=SCI,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.4,"Day 2 · Problem Solver Day",sz=16,b=True,c=ENG,a=PP_ALIGN.CENTER)
sh1=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(2.0),Inches(1.55),Inches(2.8),Inches(2.8))
sh1.fill.solid(); sh1.fill.fore_color.rgb=SCI; sh1.line.color.rgb=YELLOW; sh1.line.width=Pt(5)
tf1=tb(s,2.0,1.85,2.8,0.4,"科学家",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf1,"👩‍🔬",sz=70,a=PP_ALIGN.CENTER)
ap(tf1,"为什么？",sz=14,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(5.2),Inches(1.55),Inches(2.8),Inches(2.8))
sh2.fill.solid(); sh2.fill.fore_color.rgb=ENG; sh2.line.color.rgb=YELLOW; sh2.line.width=Pt(5)
tf2=tb(s,5.2,1.85,2.8,0.4,"工程师",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf2,"👷‍♂️",sz=70,a=PP_ALIGN.CENTER)
ap(tf2,"怎么办？",sz=14,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,1,4.65,8,0.45,"💡 谁来解决问题?  Who solves problems?",sz=16,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,1,5.15,8,0.3,"3 sessions × 50 min  ·  Hook → Mini-Lesson → Practice → Apply → Share",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"30 秒 hook:\n• 「今天三节课, 都用 5 步学习: 🔥 Hook → 📚 Mini-Lesson → 🎯 Practice → 🌱 Apply → 🎤 Share」\n• 全班按 4 队分好: 🔴 红队 / 🔵 蓝队 / 🟢 绿队 / 🟡 黄队\n• 准备好 4 色举牌 (科学家=蓝, 工程师=红) — 一会儿要投票!")

# ============================================================
# 2. TODAY'S ROADMAP — 3 sessions × 5 phases
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🗺️ 今天 3 节课  Today's 3 Sessions",NAVY)
tb(s,0.4,0.85,9.2,0.35,"每节课都有 5 步: 🔥 → 📚 → 🎯 → 🌱 → 🎤",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.28,"Each session: Hook → Mini-Lesson → Practice → Apply → Share",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sessions=[
    ("S1","🌟","故事课","Stories","上午","认识 4 位 problem solver",SCI),
    ("S2","📖","词汇课","Words","下午","收集 4 个魔法词 + 写「工」",YELLOW),
    ("S3","🧪","实验课","Experiment","下午","让脏水变干净 — 你来当工程师!",ENG),
]
for i,(tag,em,cn,en,when,desc,cl) in enumerate(sessions):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(2.95),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=cl; sh.line.fill.background()
    tb(s,x+0.05,1.80,2.85,0.45,tag,sz=18,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.25,2.85,0.85,em,sz=60,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.20,2.85,0.45,cn,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.85,0.30,en,sz=11,c=YELLOW,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.00,2.85,0.45,desc,sz=11,c=WHITE,a=PP_ALIGN.CENTER)
# Persistent leaderboard intro
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.55))
sh.fill.solid(); sh.fill.fore_color.rgb=PURPLE; sh.line.fill.background()
tb(s,0.5,4.70,9.0,0.30,"🏆 4 队全天积分 — 最高分队拿冠军徽章!",sz=14,b=True,c=WHITE)
tb(s,0.5,4.98,9.0,0.20,"4 teams · points across all 3 sessions · winners get champion badges",sz=10,c=YELLOW)
n+=1; pn(s,n)
notes(s,"1 分钟:\n• 介绍 3 节课的 flow\n• 强调 4 队全天积分 — 营造期待\n• 点出 4 队 + 每队选一个队长 (rotation)")

# ============================================================
# 3. SESSION 1 DIVIDER
# ============================================================
s=div("Session 1  上午 11:00–11:50","🌟 故事课  Problem Solver Stories  ·  50 min",SCI,"📖"); n+=1; pn(s,n)

# ============================================================
# === SESSION 1 · PHASE 1: HOOK (5 min) ===
# ============================================================
# 4. Hook slide — 以前没有电灯 + TPR
s=ns(); bg(s,CREAM); hb(s,"🌑 想象一下…  Imagine…",NAVY)
# Compact dark-room visual on left
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(0.95),Inches(3.0),Inches(2.20))
sh.fill.solid(); sh.fill.fore_color.rgb=DARK; sh.line.color.rgb=NAVY; sh.line.width=Pt(3)
tb(s,0.5,1.10,3.0,1.3,"🌑",sz=90,a=PP_ALIGN.CENTER)
tb(s,0.5,2.55,3.0,0.4,"晚上 — 一片黑!",sz=14,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,0.5,2.90,3.0,0.25,"Night — pitch dark!",sz=10,c=YELLOW,a=PP_ALIGN.CENTER)
# Right: setup + 2 questions
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.70),Inches(0.95),Inches(6.0),Inches(2.20))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=NAVY; panel.line.width=Pt(2.5)
# Setup
tb(s,3.85,1.05,5.7,0.45,"💡 以前 — 还没有电灯…",sz=15,b=True,c=NAVY)
tb(s,3.85,1.45,5.7,0.40,"一到晚上, 大家都黑乎乎的!",sz=15,b=True,c=DARK)
tb(s,3.85,1.82,5.7,0.30,"Long ago, no lightbulbs — pitch dark every night.",sz=10,c=GRAY)
# Divider
ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(3.85),Inches(2.20),Inches(5.7),Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb=LGRAY; ln.line.fill.background()
# Question 1
tb(s,3.85,2.28,5.7,0.40,"❓ 那时的生活是什么样子呢?",sz=15,b=True,c=LAB)
tb(s,3.85,2.65,5.7,0.30,"What was life like back then?",sz=10,c=GRAY)
# Question 2 — big highlight
q2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.30),Inches(9.3),Inches(0.95))
q2.fill.solid(); q2.fill.fore_color.rgb=LAB; q2.line.fill.background()
tb(s,0.5,3.40,9.1,0.45,"❓ 你知道是谁解决了这个问题吗?",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.85,9.1,0.30,"Do you know WHO solved this problem?",sz=11,c=YELLOW,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.45,"那时候人们 ___ 。 我猜是 ___ 解决的!","Back then, people ___. I guess ___ solved it!",accent=YELLOW)
n+=1; pn(s,n)
notes(s,"🔥 HOOK · 5 分钟:\n• 关掉教室一半的灯 (戏剧效果!)\n• 念 setup: 「以前 — 还没有电灯, 一到晚上, 大家都黑乎乎的!」\n• 第 1 问: 「那时的生活是什么样子?」 — 让学生想 30 秒, 让 3 个学生说\n  - 引导答案: 早睡? 用蜡烛? 用油灯? 不能看书? 不能玩?\n• 第 2 问: 「你知道是谁解决了这个问题吗?」 — 收集猜测 (爱迪生? 灯泡的人?)\n• 不揭晓答案! 「一会儿告诉你 — 现在我们先学一个 big idea」")

# ============================================================
# === SESSION 1 · PHASE 2: MINI-LESSON (10 min) ===
# ============================================================
# 5. Concept: 科学家研究什么? (comprehensive — research nature, ask WHY/HOW)
s=ns(); bg(s,CREAM); hb(s,"👩‍🔬 科学家研究什么?  Scientists ask WHY / HOW",SCI)
# Big core sentence
core=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(0.85))
core.fill.solid(); core.fill.fore_color.rgb=SCI; core.line.fill.background()
tb(s,0.55,0.92,9.0,0.36,"🌍 科学家研究自然世界:动物 · 植物 · 水 · 空气 · 身体 · 天空",sz=14,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,0.55,1.32,9.0,0.30,"Scientists study the natural world: animals, plants, water, air, body, sky",sz=10,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# 6 examples — each labeled with scientist subtype
sci_examples=[
    ("💧","为什么 水 会 脏?","→ 环境科学家  Environmental scientist"),
    ("🤒","为什么 会 生病?","→ 医学科学家  Medical scientist"),
    ("🐟","怎么 知道 鱼 怕 什么?","→ 动物科学家  Animal behaviorist"),
    ("⭐","为什么 星星 会 亮?","→ 天文学家  Astronomer"),
    ("🦖","为什么 恐龙 不见了?","→ 古生物学家  Paleontologist"),
    ("☀️","太阳 怎么 给 我们 温暖?","→ 物理学家  Physicist"),
]
for i,(em,q_cn,subtype) in enumerate(sci_examples):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.85+row*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=SCI; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.18,0.50,0.5,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.65,y+0.06,2.25,0.32,q_cn,sz=11,b=True,c=DARK)
    tb(s,x+0.65,y+0.36,2.25,0.42,subtype,sz=8,b=True,c=SCI)
# Bottom: 科学家研究 summary box (topics + key questions)
hint=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.45))
hint.fill.solid(); hint.fill.fore_color.rgb=WARM; hint.line.color.rgb=SCI; hint.line.width=Pt(2.5)
tb(s,0.55,3.92,9.0,0.30,"🔬 科学家研究: Scientists study:",sz=12,b=True,c=SCI)
# Topic icons row
topics=[("🌱","植物"),("🐟","动物"),("💧","水"),("🌬","空气"),("🤒","身体"),("⭐","天空"),("🦖","以前的生命")]
tx=0.55
for em,cn in topics:
    tb(s,tx,4.22,1.30,0.30,em,sz=20,a=PP_ALIGN.LEFT)
    tb(s,tx+0.30,4.30,1.0,0.24,cn,sz=10,b=True,c=DARK)
    tx+=1.30
# Key questions
tb(s,0.55,4.65,9.0,0.30,"❓ 想知道: 为什么? 怎么知道?",sz=14,b=True,c=SCI,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.24,"Want to know: WHY? HOW DO WE KNOW?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"📚 MINI 1 · 4-5 分钟:\n• 关键讲法: 「科学家不是随便问问题。科学家会观察自然, 然后问: 为什么会这样? 我们怎么知道?」\n• 念中间大句子 — 让全班跟读: 「科学家研究自然世界 — 动物、植物、水、空气、身体、天空」\n• 念 6 个例子 — 强调每个问题都对应一种 SCIENTIST:\n  - 水脏 → 环境科学家\n  - 生病 → 医学科学家\n  - 鱼怕什么 → 动物科学家 (观察行为)\n  - 星星亮 → 天文学家\n  - 恐龙 → 古生物学家\n  - 太阳温暖 → 物理学家\n• 关键: 用「怎么」也是科学家 — 只要是研究自然的, 就是科学家!\n• 课堂互动:\n  - 「为什么会生病? 是科学家的问题吗?」 → 是!\n  - 「怎么做一个机器人? 是科学家还是工程师?」 → 工程师!\n• 总结: 科学家 找答案, 工程师 做办法。")

# 6. Concept: 工程师解决什么? (comprehensive — solve life problems)
s=ns(); bg(s,CREAM); hb(s,"👷‍♂️ 工程师解决什么?  Engineers ask HOW can we solve it?",ENG)
# Big core sentence
core=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(0.85))
core.fill.solid(); core.fill.fore_color.rgb=ENG; core.line.fill.background()
tb(s,0.55,0.92,9.0,0.36,"🔧 工程师设计东西, 解决生活问题:桥 · 房子 · 机器 · 电脑 · 机器人 · 交通灯",sz=13,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,0.55,1.32,9.0,0.30,"Engineers design things to solve life problems: bridges, houses, machines, computers, robots, lights",sz=9,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# 6 examples — each labeled with engineer subtype
eng_examples=[
    ("🌊","怎么 安全 过河?","→ 土木工程师  Civil engineer"),
    ("🏠","房子 怎么 不倒?","→ 结构工程师  Structural engineer"),
    ("💡","灯 怎么 亮 起来?","→ 电气工程师  Electrical engineer"),
    ("💻","游戏 / App 怎么 动?","→ 软件工程师  Software engineer"),
    ("🤖","机器人 怎么 帮人?","→ 机器人工程师  Robotics engineer"),
    ("🚦","马路 太乱 怎么办?","→ 交通工程师  Traffic engineer"),
]
for i,(em,q_cn,subtype) in enumerate(eng_examples):
    col=i%3; row=i//3
    x=0.4+col*3.10; y=1.85+row*0.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(0.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=ENG; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.18,0.50,0.5,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.65,y+0.06,2.25,0.32,q_cn,sz=11,b=True,c=DARK)
    tb(s,x+0.65,y+0.36,2.25,0.42,subtype,sz=8,b=True,c=ENG)
# Bottom: 工程师解决 summary box
hint=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.45))
hint.fill.solid(); hint.fill.fore_color.rgb=WARM; hint.line.color.rgb=ENG; hint.line.width=Pt(2.5)
tb(s,0.55,3.92,9.0,0.30,"🛠️ 工程师解决: Engineers solve:",sz=12,b=True,c=ENG)
# Topic icons row
topics=[("🌉","过河"),("🏠","房子"),("💡","用电"),("💻","电脑"),("🤖","机器人"),("🚦","交通"),("📱","通讯")]
tx=0.55
for em,cn in topics:
    tb(s,tx,4.22,1.30,0.30,em,sz=20,a=PP_ALIGN.LEFT)
    tb(s,tx+0.30,4.30,1.0,0.24,cn,sz=10,b=True,c=DARK)
    tx+=1.30
# Key questions
tb(s,0.55,4.65,9.0,0.30,"❓ 想知道: 怎么办? 怎么做? 怎么更好?",sz=14,b=True,c=ENG,a=PP_ALIGN.CENTER)
tb(s,0.55,4.95,9.0,0.24,"Want to know: HOW? HOW to make? HOW to make it better?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"📚 MINI 2 · 4-5 分钟:\n• 关键讲法: 「工程师不是只会做东西。工程师先看到一个问题, 然后想办法、设计、制作、测试, 再改进。」\n• 念中间大句子 — 让全班跟读: 「工程师设计东西, 解决生活问题 — 桥、房子、机器、电脑、机器人、交通灯」\n• 念 6 个例子 — 强调每个问题对应一种 ENGINEER:\n  - 过河 → 土木工程师 (设计桥)\n  - 房子不倒 → 结构工程师\n  - 灯亮 → 电气工程师 (电路)\n  - 游戏/App → 软件工程师 (写程序)\n  - 机器人 → 机器人工程师\n  - 交通 → 交通工程师\n• 课堂互动 (对比 SCIENTIST vs ENGINEER):\n  - 「水为什么会脏? 是科学家还是工程师?」 → 科学家! (找原因)\n  - 「怎么让水变干净? 是科学家还是工程师?」 → 工程师! (做办法)\n• 总结: 科学家 找答案, 工程师 做办法!")

# ============================================================
# 9.5 · 8 QUESTION BANK — mixed colors, varying difficulty (some tricky!)
# Card colors are NON-correlated with answer — kids can't infer from color
# Some questions break the "为什么=科 / 怎么=工" pattern
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"📝 题库 · 8 个问题  Question Bank",PURPLE)
group_label(s)
score_badge(s)
tb(s,0.4,1.20,9.2,0.40,"⚠️ 小心 trick! 「怎么」 不一定 = 工程师, 「为什么」 不一定 = 科学家!",sz=14,b=True,c=LAB,a=PP_ALIGN.CENTER)
# 8 mixed-difficulty questions; card colors deliberately NOT grouped by category
questions=[
    ("🌙","为什么 月亮 有时圆 有时弯?","Why is the moon round/curved?","科",SCI),
    ("🏠","怎么 让 房子 不 怕 地震?","Earthquake-proof house?","工",PURPLE),
    ("🐟","怎么 知道 鱼 怕 什么?","How do we know what fish fear?","科",GREEN),
    ("✈️","为什么 飞机 会 飞?","Why can planes fly?","科",EDISON),
    ("👴","怎么 帮 老爷爷 上楼?","How to help grandpa go up?","工",MORGAN),
    ("☀️","太阳 怎么 给 我们 温暖?","How does sun warm us?","科",WRIGHT),
    ("📺","没电了 怎么 看 电视?","No power — how to watch TV?","工",GOODALL),
    ("💭","为什么 我们 会 做梦?","Why do we dream?","科",ENG),
]
for i,(em,q_cn,q_en,ans,cl) in enumerate(questions):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.70+row*0.83
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.78))
    sh.fill.solid(); sh.fill.fore_color.rgb=cl; sh.line.fill.background()
    # Number badge (white circle)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.18),Inches(0.42),Inches(0.42))
    nb.fill.solid(); nb.fill.fore_color.rgb=WHITE; nb.line.fill.background()
    tb(s,x+0.10,y+0.22,0.42,0.34,str(i+1),sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    # Emoji
    tb(s,x+0.55,y+0.18,0.50,0.42,em,sz=20,a=PP_ALIGN.CENTER)
    # Question Chinese + English
    tb(s,x+1.10,y+0.06,2.85,0.32,q_cn,sz=12,b=True,c=WHITE)
    tb(s,x+1.10,y+0.42,2.85,0.30,q_en,sz=8,c=YELLOW)
    # Answer chip — small colored circle (科=teal, 工=rust) revealing category
    ans_color = SCI if ans=="科" else ENG
    ac=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+4.05),Inches(y+0.18),Inches(0.42),Inches(0.42))
    ac.fill.solid(); ac.fill.fore_color.rgb=ans_color; ac.line.color.rgb=YELLOW; ac.line.width=Pt(2)
    tb(s,x+4.05,y+0.22,0.42,0.34,ans,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# Bottom: rules
rules=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.05),Inches(9.2),Inches(0.30))
rules.fill.solid(); rules.fill.fore_color.rgb=YELLOW; rules.line.fill.background()
tb(s,0.5,5.09,9.0,0.22,"🏆 第一指对 +2 分  ·  😊 答案在右边小圆圈 — 别 cheat!",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5 分钟 PRACTICE — 8 题, 难度递增:\n• Tricky 题 (4): 「为什么飞机会飞」 → 科学家 (空气动力学), 不是工程师!\n• Tricky 题 (3): 「怎么知道鱼怕什么」 → 科学家 (研究方法), 不是工程师!\n• Tricky 题 (6): 「太阳怎么给我们温暖」 → 用「怎么」 但是问 自然现象 → 科学家\n• 顺序可以打乱\n• 老师先盖住答案 (一手挡), 让各队猜后再揭晓\n• 答错的队不扣分 — 鼓励试错\n• 收尾: 「你看 — '怎么' 不一定是工程师, '为什么' 不一定是科学家! 要看 — 是想 understand 还是 design」")

# ============================================================
# === SESSION 1 · PHASE 3: ACTIVE PRACTICE (20 min) ===
# 4 problem solvers × 5 min each (story 3 min + Q&A 2 min)
# Story FIRST (kids haven't learned), then interactive Q&A
# ============================================================
# --- Edison: VIDEO ---
s=video_slide("💡","爱迪生","Thomas Edison",EDISON,
    "搜: 'Thomas Edison kids' / BrainPOP Jr. 'Thomas Edison'  ·  绘本: 「The Boy Who Loved Math」 系列",
    "Search 'Thomas Edison for kids' (3-5 min)",
    "他失败了几次? 怎么不放弃?","How many failures? Why didn't he give up?")
n+=1; pn(s,n)
notes(s,"先看视频 (1-2 分钟):\n• 老师课前找好链接\n• 推荐: BrainPOP Jr. 'Thomas Edison' (英文) / 'Thomas Edison 中文动画'\n• 看的时候让学生数: 他失败了多少次?\n• 看完直接翻下一页 — 一起回答 3 个问题")

# --- Edison: STORY + Q&A ---
s=story_slide("💡","爱迪生","Thomas Edison","1847–1931","🇺🇸 美国",EDISON,
    [("❓","问题是什么?","What was the problem?",
       "100 年前晚上一片黑 — 不能看书, 不能玩.",
       "100 yrs ago: pitch-dark nights — couldn't read or play.",
       "你看完视频, 你说!"),
     ("❓","他做了什么?","What did he do?",
       "他试了 1000 多次! 每次失败都不放弃, 继续试.",
       "He tried 1000+ times! Never gave up after failures.",
       "他失败了几次? 你猜!"),
     ("❓","结果是什么?","What was the result?",
       "灯泡亮了! 现在每个房间都有灯, 晚上也很亮.",
       "The bulb lit up! Now every room has light.",
       "啪! — 一起拧灯泡!")],
    "🎙️ 他还发明了留声机 — 第一次能录音乐!","He also invented the phonograph — first music recorder!",
    video_cn="搜: 'Thomas Edison kids' / BrainPOP Jr. 'Thomas Edison'  ·  绘本: 「The Boy Who Loved Math」 系列",
    video_en="Search 'Thomas Edison for kids' (3-5 min)")
n+=1; pn(s,n)
notes(s,"5 分钟:\n🎥 1-2 分钟看视频 (老师课前找好链接)\n  推荐: BrainPOP Jr. 'Thomas Edison' (英文) / 'Thomas Edison 中文动画'\n  绘本备选: 'Edison's Concrete Piano' / 'A Wizard from the Start' (节选页)\n📋 3 分钟答 3 题 (轮流上台答):\n• Q1 问题是什么? → 让学生说看到的, 不直接给答案\n• Q2 他做了什么? → 强调 失败 1000 次\n• Q3 结果是什么? → 全班拧灯泡 「啪!」\n• 答对每题 → 队 +1 分")

# --- Jane Goodall: VIDEO ---
s=video_slide("🐒","简·古道尔","Jane Goodall",GOODALL,
    "绘本: 「Me … Jane」 by Patrick McDonnell  ·  搜: 'Jane Goodall for kids' / 'Roots & Shoots Jane'",
    "Picture book: 'Me…Jane' / Search 'Jane Goodall kids documentary'",
    "她在哪里? 她在做什么?","Where is she? What is she doing?")
n+=1; pn(s,n)
notes(s,"先看视频 (1-2 分钟):\n• 推荐: 'Me…Jane' picture book read-aloud (YouTube 有), 或 National Geographic Kids 'Jane Goodall'\n• 让学生看: 她一个人观察了多久?\n• 看完直接翻下一页")

# --- Jane Goodall: STORY + Q&A ---
s=story_slide("🐒","简·古道尔","Jane Goodall","1934– 至今","🇬🇧 英国",GOODALL,
    [("❓","简提出的问题是什么?","What was Jane's question?",
       "黑猩猩在森林里怎么生活? 它们想什么?",
       "How do chimps live in the forest? What do they think?",
       "你看完视频/绘本, 说说看!"),
     ("❓","她做了什么?","What did she do?",
       "她一个人去 非洲森林, 安静地观察黑猩猩 — 一年又一年.",
       "Lived alone in African forest, silently watched chimps for years.",
       "你能安静 1 分钟吗?"),
     ("❓","结果是什么?","What was the result?",
       "她发现 — 黑猩猩用树枝钓白蚁! 它们也会用工具!",
       "She found chimps use twigs as tools — like humans!",
       "太神奇了! 你猜还有什么发现?")],
    "💚 她给每只黑猩猩起了名字 — 像家人一样!","She named every chimp — like family!",
    video_cn="绘本: 「Me ... Jane」 by Patrick McDonnell  ·  搜: 'Jane Goodall for kids' / 'Roots & Shoots Jane'",
    video_en="Picture book: 'Me…Jane' / Search 'Jane Goodall kids documentary'")
n+=1; pn(s,n)
notes(s,"5 分钟 (这是今天唯一的 SCIENTIST!):\n🎥 1-2 分钟 — 推荐: 'Me…Jane' picture book read-aloud (YouTube 有), 或 National Geographic Kids 'Jane Goodall'\n📋 3 分钟答 3 题:\n• Q1 她的问题: 黑猩猩怎么生活/想什么? (重点: 她想 understand)\n• Q2 她做了什么: 强调 一个人, 安静, 观察 — 不打扰\n• Q3 结果: 黑猩猩用工具! 全班「安静挑战」 30 秒\n• 「这就是 SCIENTIST — 想 understand 自然」")

# --- Wright Brothers: VIDEO ---
s=video_slide("✈️","莱特兄弟","Wright Brothers",WRIGHT,
    "绘本: 「My Brothers' Flying Machine」 by Jane Yolen  ·  搜: 'Wright Brothers first flight kids'",
    "Picture book: 'My Brothers' Flying Machine' / Search 'Wright Brothers for kids'",
    "他们试了多少次? 第一次飞了多久?","How many tries? How long was the first flight?")
n+=1; pn(s,n)
notes(s,"先看视频 (1-2 分钟):\n• 推荐: 'Wright Brothers First Flight in 12 Seconds' YouTube\n• 绘本: 'My Brothers' Flying Machine' (Jane Yolen)\n• 让学生看: 他们摔了几次? 第一次飞了几秒?\n• 看完翻下一页一起答")

# --- Wright Brothers: STORY + Q&A ---
s=story_slide("✈️","莱特兄弟","Wright Brothers","1903 第一次飞","🇺🇸 美国",WRIGHT,
    [("❓","他们想解决什么?","What did they want to solve?",
       "几千年来人想飞 — 可是不会! 兄弟想:「能不能让人飞?」",
       "For 1000s of years humans couldn't fly — they asked: Can we?",
       "你想飞吗? 飞到哪里?"),
     ("❓","他们做了什么?","What did they do?",
       "在沙滩上试了几百次! 哥哥弟弟 一起做飞机, 摔了又改, 改了又摔.",
       "Hundreds of beach trials. Built, crashed, fixed, repeated.",
       "你猜他们摔了几次?"),
     ("❓","结果是什么?","What was the result?",
       "1903 年成功! 第一次飞 12 秒, 36 米 — 短, 但是开始了!",
       "1903: First flight! 12 sec, 36 m — short, but a start!",
       "现在飞机能飞多远?")],
    "🏛️ 第一架飞机叫 Flyer — 现在还在博物馆里!","First plane was called Flyer — still in a museum today!",
    video_cn="绘本: 「My Brothers' Flying Machine」 by Jane Yolen  ·  搜: 'Wright Brothers first flight kids'",
    video_en="Picture book: 'My Brothers' Flying Machine' / Search 'Wright Brothers for kids'")
n+=1; pn(s,n)
notes(s,"5 分钟:\n🎥 1-2 分钟 — 推荐: 'Wright Brothers First Flight in 12 Seconds' YouTube\n  绘本: 'My Brothers' Flying Machine' (Jane Yolen) — 弟弟 Katherine 视角\n📋 3 分钟答:\n• Q1 想解决什么: 强调「为什么人不能飞?」 — 但他们更想问 「怎么让人飞?」 → 工程师!\n• Q2 做了什么: 几百次失败 — 跟爱迪生一样, 不放弃\n• Q3 结果: 12 秒 — 让学生比较现在 (10 多小时!) — 工程进步\n• TPR bonus: 双臂展开, 教室转 1 圈再坐下")

# --- Garrett Morgan: VIDEO ---
s=video_slide("🚦","加勒特·摩根","Garrett Morgan",MORGAN,
    "搜: 'Garrett Morgan inventor for kids'  ·  绘本: 「Garrett Morgan, Inventor Extraordinaire」",
    "Picture book: 'Garrett Morgan, Inventor Extraordinaire' / Search 'Garrett Morgan kids'",
    "他看到什么问题? 他做了什么?","What problem did he see? What did he do?")
n+=1; pn(s,n)
notes(s,"先看视频 (1-2 分钟):\n• 推荐: PBS 'Garrett Morgan' kids segment, 或绘本 read-aloud\n• Black History 教学好机会\n• 让学生看: 他看到的问题是什么? 他怎么解决?\n• 看完翻下一页")

# --- Garrett Morgan: STORY + Q&A ---
s=story_slide("🚦","加勒特·摩根","Garrett Morgan","1877–1963","🇺🇸 美国",MORGAN,
    [("❓","摩根看到什么问题?","What problem did Morgan see?",
       "100 年前马路没有红绿灯! 车 太 多, 经常撞 — 他亲眼看到一次大车祸.",
       "No traffic lights 100 yrs ago — many crashes. He saw a bad one.",
       "你过马路怕吗? 为什么?"),
     ("❓","他做了什么?","What did he do?",
       "他设计了 红 黄 绿 三色 信号灯 — 让车 停, 慢, 走 都有 信号.",
       "Designed 3-color signals so cars know stop / slow / go.",
       "你猜 — 红是 停 还是 走?"),
     ("❓","结果是什么?","What was the result?",
       "马路安全多了! 现在 全世界的 红绿灯 都从他的设计来.",
       "Roads got safer! Today's traffic lights all from his design.",
       "全班齐喊: 红停! 绿走! 黄慢!")],
    "😷 他还发明了防毒面具 — 救了火灾里的人!","He also invented the gas mask — saved lives in fires!",
    video_cn="搜: 'Garrett Morgan inventor for kids'  ·  绘本: 「Garrett Morgan, Inventor Extraordinaire」",
    video_en="Picture book: 'Garrett Morgan, Inventor Extraordinaire' / Search 'Garrett Morgan kids'")
n+=1; pn(s,n)
notes(s,"5 分钟:\n🎥 1-2 分钟 — 推荐: PBS 'Garrett Morgan' kids segment, 或绘本 read-aloud\n  Black History 教学好机会 — 强调他作为非裔美国人的成就\n📋 3 分钟答:\n• Q1 看到什么问题: 一个真实的车祸触发了他\n• Q2 做了什么: 设计 + 申请专利 — 工程师\n• Q3 结果: 红绿灯救人 — 让学生想现在路上有多少红绿灯\n• 红绿灯游戏: 红→站住, 绿→跑步, 黄→拍手 (1 分钟 TPR)")

# ============================================================
# === SESSION 1 · PHASE 4: APPLY (10 min) ===
# ============================================================
# Apply 1 — 谁来解决? Whole-class match (no groups, full questions)
s=ns(); bg(s,CREAM); hb(s,"🎮 谁来解决?  Who Solved This Problem?",LAB)
tb(s,0.4,0.85,9.2,0.40,"老师念完整的问题 — 全班一起想:这是谁解决的?",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Teacher reads the full problem — class thinks: who solved it?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
matches=[
    ("🌑","100 年前晚上一片黑 — 不能看书 不能玩。","100 yrs ago: pitch-dark nights — couldn't read or play.","💡","爱迪生 (灯泡)","Edison (lightbulb)",EDISON),
    ("🐒","想 understand 黑猩猩怎么生活、怎么用工具。","Wanted to understand how chimps live and use tools.","🌿","简·古道尔 (科学家)","Goodall (scientist)",GOODALL),
    ("✈️","几千年来人想飞 — 但是没有翅膀, 不会飞。","For 1000s of years humans wanted to fly — but couldn't.","🛩️","莱特兄弟 (飞机)","Wright Brothers (airplane)",WRIGHT),
    ("🚦","100 年前马路没有红绿灯 — 经常出车祸。","100 yrs ago: no traffic lights — many crashes.","🚥","加勒特·摩根 (红绿灯)","Garrett Morgan (traffic light)",MORGAN),
]
for i,(em1,p_cn,p_en,em2,a_cn,a_en,cl) in enumerate(matches):
    y=1.65+i*0.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.78))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2)
    # Problem (left half)
    tb(s,0.55,y+0.08,0.5,0.35,em1,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.10,y+0.06,4.0,0.36,p_cn,sz=12,b=True,c=DARK)
    tb(s,1.10,y+0.42,4.0,0.30,p_en,sz=9,c=GRAY)
    # Arrow
    tb(s,5.20,y+0.18,0.5,0.40,"→",sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    # Solver (right half)
    tb(s,5.85,y+0.08,0.5,0.35,em2,sz=22,a=PP_ALIGN.CENTER)
    tb(s,6.45,y+0.06,3.10,0.36,a_cn,sz=13,b=True,c=cl)
    tb(s,6.45,y+0.42,3.10,0.30,a_en,sz=9,c=GRAY)
# Bottom prompt — whole-class sentence
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(5.10),Inches(9.2),Inches(0.45))
prompt.fill.solid(); prompt.fill.fore_color.rgb=WARM; prompt.line.color.rgb=LAB; prompt.line.width=Pt(1.5)
tb(s,0.55,5.16,9.0,0.32,"💬 全班一起说: ___ 解决了 ___ 的问题。 / ___ solved the ___ problem.",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🌱 APPLY 1 · 5 分钟 — 全班一起做 (不分组, 不计分):\n• 老师念完整的问题 (左侧 Chinese), 让全班一起想 — 这是谁解决的?\n• 完整问题让小孩子也能听懂上下文\n• 让 1-2 个学生说出名字, 然后老师揭示右侧答案\n• 全班齐声: 「___ 解决了 ___ 的问题」 — 用句型说出每一对\n• 4 对都做完, 重点回顾: 谁是 SCIENTIST? 谁是 ENGINEER?\n• Goodall 是唯一的 SCIENTIST — Edison/Wright/Morgan 都是 ENGINEER")

# Apply 2 — 我想解决___ 分组讨论 (group brainstorm with multiple example prompts)
s=ns(); bg(s,CREAM); hb(s,"💭 我想解决…  What I Want to Fix · 分组讨论",GREEN)
tb(s,0.4,0.85,9.2,0.36,"想一想 — 你想解决什么问题? 看看下面的例子!",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.22,9.2,0.26,"Think — what problem do YOU want to fix? See examples below!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 4 categories with multiple example prompts each
prompts=[
    ("🏫","学校 School",["作业太多","食堂的菜不好吃","操场太小","没有空调"],GREEN),
    ("🏠","家 Home",["垃圾多","房间乱","弟弟妹妹爱哭","网络太慢"],EDISON),
    ("🌍","世界 World",["海里有塑料","雾霾天气","空气脏","水不够喝"],SCI),
    ("🐾","动物 Animals",["流浪猫狗","动物受伤","濒危动物","海龟吃塑料"],BROWN),
]
for i,(em,cn_label,examples,cl) in enumerate(prompts):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    # Emoji + label
    tb(s,x+0.10,y+0.10,0.7,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.13,3.6,0.40,cn_label,sz=14,b=True,c=cl)
    # Example list
    for j,ex in enumerate(examples):
        tb(s,x+0.20,y+0.62+j*0.22,4.20,0.22,f"·  {ex}",sz=10,c=DARK)
# Bottom: discussion + sentence frame
disc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.70))
disc.fill.solid(); disc.fill.fore_color.rgb=GREEN; disc.line.fill.background()
tb(s,0.55,4.92,9.0,0.30,"👥 3-4 人一组讨论 — 选一个例子或自己想一个! 然后分享给全班",sz=12,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.30,"💬 我想解决 ___ 的问题, 我要 ___ 。",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🌱 APPLY 2 · 6-8 分钟 — 分组讨论:\n• Step 1 (1 分钟): 个人看例子, 想一想\n• Step 2 (3 分钟): 3-4 人一组讨论\n  - 每人选一个例子 (或自己想一个)\n  - 用句型说: 「我想解决 ___ 的问题, 我要 ___ 。」\n  - 组员互相听, 给建议\n• Step 3 (3 分钟): 每组选 1 位代表上台分享\n• 不评判好坏 — 鼓励每个想法\n• 老师可以补充: 每个想法都是未来 SCIENTIST/ENGINEER 的开始!")

# ============================================================
# === SESSION 1 · PHASE 5: SHARE & CLOSE (5 min) ===
# ============================================================

# Combined Summary + Transition (no separate phase marker — quick close)
s=ns(); bg(s,CREAM); hb(s,"🎤 总结 + 下一步  Summary + Next!",PH_CLOSE)
score_badge(s)
# Quick recap
tb(s,0.4,0.85,9.2,0.4,"🧭 今早学了什么?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.40),Inches(4.4),Inches(1.85))
left.fill.solid(); left.fill.fore_color.rgb=SCI; left.line.fill.background()
tb(s,0.5,1.50,4.4,0.5,"👩‍🔬 科学家",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.00,4.4,0.55,"为什么?",sz=24,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,0.5,2.55,4.4,0.30,"WHY?  ·  简·古道尔 🐒",sz=12,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,2.85,4.4,0.30,"研究、观察、找答案",sz=11,c=YELLOW,a=PP_ALIGN.CENTER)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.40),Inches(4.4),Inches(1.85))
right.fill.solid(); right.fill.fore_color.rgb=ENG; right.line.fill.background()
tb(s,5.10,1.50,4.4,0.5,"👷‍♂️ 工程师",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,5.10,2.00,4.4,0.55,"怎么办?",sz=24,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,5.10,2.55,4.4,0.30,"HOW?  ·  💡✈️🚦",sz=12,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,5.10,2.85,4.4,0.30,"设计、建造、改进",sz=11,c=YELLOW,a=PP_ALIGN.CENTER)
# Transition strip
trans=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.40),Inches(9.0),Inches(1.45))
trans.fill.solid(); trans.fill.fore_color.rgb=LAB; trans.line.color.rgb=YELLOW; trans.line.width=Pt(3)
tb(s,0.5,3.50,9.0,0.45,"💧 下一个挑战 → 水很脏! 怎么办?",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.95,9.0,0.30,"Next challenge → Dirty water! How to fix it?",sz=11,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,0.5,4.30,9.0,0.40,"🧪 下午 — 你们来当工程师!",sz=15,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎤 SHARE & CLOSE · 5 分钟:\n• 1 分钟: 全班齐喊 — 「科学家! 为什么?」 「工程师! 怎么办?」 (各 3 次, 加手势)\n• 1 分钟: 每队代表说 1 句 — 「我们队学到了 ___ 」 (轮流上台)\n• 1 分钟: 公布 Session 1 积分 — 暂时领先的队\n• 2 分钟: Tease 下午实验 — 「水很脏! 怎么办? — 一会儿告诉你!」 (留悬念)")

# ============================================================
# 25. SESSION 2 DIVIDER (Day 1-style: Review + Read + Write)
# ============================================================
s=div("Session 2  下午 1:00–1:50","📚 复习 + 我会认 + 我会写  Review · Read · Write",YELLOW,"📖"); n+=1; pn(s,n)

# ============================================================
# 26. REVIEW — Session 1 recap (4 cards covering scientist/engineer + 4 figures)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🔄 复习  Review · Session 1",NAVY)
tb(s,0.4,0.85,9.2,0.40,"早上学了什么？  What did we learn this morning?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Quick recap before we read & write!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
recap=[
    ("👩‍🔬","科学家 = 为什么?","Scientist asks WHY","简·古道尔 · 观察黑猩猩",SCI),
    ("👷‍♂️","工程师 = 怎么办?","Engineer asks HOW","爱迪生 · 莱特兄弟 · 摩根",ENG),
    ("💡","他们都解决问题","They all solve problems","失败也不放弃 — 再试!",EDISON),
    ("💭","我也想…","I also want to…","「我想 ___ , 因为 ___ 。」",GREEN),
]
for i,(em,cn,en,detail,c) in enumerate(recap):
    col=i%2; row=i//2
    x=0.4+col*4.65
    y=1.65+row*1.45
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.30))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.12,0.7,0.65,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.10,3.6,0.40,cn,sz=15,b=True,c=c)
    tb(s,x+0.85,y+0.45,3.6,0.30,en,sz=10,c=GRAY)
    tb(s,x+0.15,y+0.85,4.30,0.40,detail,sz=11,b=True,c=DARK)
sentence_frame_bar(s,4.65,
    "今天早上我学了 ___ 。",
    "This morning I learned about ___.")
n+=1; pn(s,n)
notes(s,"5 分钟 review (don't tell — let students recall):\n• 「早上学了哪 4 位 problem solver?」抢答\n• 复习: 科学家 = 为什么? 工程师 = 怎么办?\n• 强调: 简·古道尔 是 唯一的 scientist; 其他 3 位是 engineer\n• 抽 1-2 个学生说 「我想 ___ 因为 ___」")

# ============================================================
# 27-32. 我会认 — 6 vocabulary words, ONE SLIDE PER WORD (Day 1 pattern)
# ============================================================
read_words=[
    ("👩‍🔬","科学家","kē xué jiā","Scientist",SCI,
        "简·古道尔 是 科学家 — 她研究黑猩猩。",
        "📷 科学家 / 实验室 / 显微镜"),
    ("👷‍♂️","工程师","gōng chéng shī","Engineer",ENG,
        "爱迪生 是 工程师 — 他发明了电灯。",
        "📷 工程师 / 安全帽 / 蓝图"),
    ("🧪","实验","shí yàn","Experiment",LAB,
        "我们做 实验 — 让脏水变干净!",
        "📷 烧杯 / 试管 / 实验"),
    ("💡","发明","fā míng","Invention",EDISON,
        "电灯 是 一个 大 发明!",
        "📷 灯泡 / 发明 / 留声机"),
    ("👀","观察","guān chá","Observe",GOODALL,
        "简 在森林里 观察 黑猩猩。",
        "📷 放大镜 / 观察 / 眼睛"),
    ("🌍","环境","huán jìng","Environment",GREEN,
        "我们要保护 环境 — 让地球更干净。",
        "📷 地球 / 树林 / 河流"),
]
for em,cn,py,en,c,sent,img_label in read_words:
    s=ns(); bg(s,CREAM); hb(s,f"👀 我会认 · {cn}  I Can Read",c)
    # Left: big character card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.fill.background()
    tb(s,0.5,1.10,4.3,1.4,cn,sz=64,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.40,4.3,0.4,f"{py}  {en}",sz=18,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读!  Read after me!",sz=14,c=c,a=PP_ALIGN.CENTER)
    # Right: image placeholder box
    ib_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.3),Inches(1.0),Inches(4.4),Inches(2.5))
    ib_box.fill.solid(); ib_box.fill.fore_color.rgb=IMGBG; ib_box.line.fill.background()
    tb(s,5.3,2.05,4.4,0.4,img_label,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    # Bottom: example sentence
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=c; sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=c)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    n+=1; pn(s,n)
    notes(s,f"3 分钟 — {cn}:\n• 老师指字, 全班齐读 3 遍 (慢→快→唱)\n• 看图: 「这是 ___ , 在做什么?」\n• 读例句, 跟读\n• 抽 1-2 个学生用 {cn} 造一个新句子\n• 写到黑板上 — 让学生 trace 一下 (空中写)")

# ============================================================
# 33. 我会写 · 实验 (stroke order)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 实验  I Can Write · Experiment",LAB)
tb(s,0.4,0.85,9.2,0.40,"练一练 — 写「实验」!",sz=22,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 实验 (Experiment)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
char1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(2.30),Inches(2.85))
char1.fill.solid(); char1.fill.fore_color.rgb=WHITE
char1.line.color.rgb=LAB; char1.line.width=Pt(3)
tb(s,0.4,1.95,2.30,1.95,"实",sz=130,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,2.30,0.30,"shí (real / true)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
char2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.85),Inches(1.65),Inches(2.30),Inches(2.85))
char2.fill.solid(); char2.fill.fore_color.rgb=WHITE
char2.line.color.rgb=LAB; char2.line.width=Pt(3)
tb(s,2.85,1.95,2.30,1.95,"验",sz=130,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,2.85,4.10,2.30,0.30,"yàn (test / examine)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
panel.line.color.rgb=LAB; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=LAB; head.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写  How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"1️⃣ 「实」 — 8 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  上面「宀」(房顶), 下面「头」",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"2️⃣ 「验」 — 10 笔",sz=14,b=True,c=DARK)
tb(s,5.45,3.40,4.10,0.30,"  左边「马」, 右边「佥」",sz=10,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格练 3 遍",sz=12,b=True,c=LAB)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「实验」! 我做了 ___ 实验。",
    "I can write 实验! I did a ___ experiment.")
n+=1; pn(s,n)
notes(s,"5-6 分钟:\n• 演示笔顺, 学生跟写 (空中写)\n• 田字格练 3 遍\n• 提示: 「实」上面像房顶 (宀); 「验」左边是马 — 「马的考试 = 实验」\n• 让学生说今天下午要做什么实验 (filter 脏水)")

# ============================================================
# 34. 我会写 · 发明 (stroke order)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 发明  I Can Write · Invention",EDISON)
tb(s,0.4,0.85,9.2,0.40,"练一练 — 写「发明」!",sz=22,b=True,c=EDISON,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 发明 (Invention)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
char3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(2.30),Inches(2.85))
char3.fill.solid(); char3.fill.fore_color.rgb=WHITE
char3.line.color.rgb=EDISON; char3.line.width=Pt(3)
tb(s,0.4,1.95,2.30,1.95,"发",sz=130,b=True,c=EDISON,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,2.30,0.30,"fā (send out)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
char4=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.85),Inches(1.65),Inches(2.30),Inches(2.85))
char4.fill.solid(); char4.fill.fore_color.rgb=WHITE
char4.line.color.rgb=EDISON; char4.line.width=Pt(3)
tb(s,2.85,1.95,2.30,1.95,"明",sz=130,b=True,c=EDISON,a=PP_ALIGN.CENTER)
tb(s,2.85,4.10,2.30,0.30,"míng (bright)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
panel2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel2.fill.solid(); panel2.fill.fore_color.rgb=WHITE
panel2.line.color.rgb=EDISON; panel2.line.width=Pt(2.5)
head2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head2.fill.solid(); head2.fill.fore_color.rgb=EDISON; head2.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写  How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"1️⃣ 「发」 — 5 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  从左上往下, 像「友」+ 短撇",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"2️⃣ 「明」 — 8 笔",sz=14,b=True,c=DARK)
tb(s,5.45,3.40,4.10,0.30,"  日 (太阳) + 月 (月亮) = 明亮!",sz=10,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格练 3 遍",sz=12,b=True,c=EDISON)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「发明」! 我想 发明 ___ 。",
    "I can write 发明! I want to invent ___.")
n+=1; pn(s,n)
notes(s,"5-6 分钟:\n• 演示笔顺\n• 提示: 「明」 = 日 (太阳) + 月 (月亮) = 明亮! 这是字的「故事」, 学生很容易记\n• 让学生说想发明什么 — 「我想发明 ___ !」 (轮流上台)")

# ============================================================
# 35. 我会写 · 观察 (stroke order)
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 观察  I Can Write · Observe",GOODALL)
tb(s,0.4,0.85,9.2,0.40,"练一练 — 写「观察」!",sz=22,b=True,c=GOODALL,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 观察 (Observe)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
char5=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(2.30),Inches(2.85))
char5.fill.solid(); char5.fill.fore_color.rgb=WHITE
char5.line.color.rgb=GOODALL; char5.line.width=Pt(3)
tb(s,0.4,1.95,2.30,1.95,"观",sz=130,b=True,c=GOODALL,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,2.30,0.30,"guān (look at)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
char6=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.85),Inches(1.65),Inches(2.30),Inches(2.85))
char6.fill.solid(); char6.fill.fore_color.rgb=WHITE
char6.line.color.rgb=GOODALL; char6.line.width=Pt(3)
tb(s,2.85,1.95,2.30,1.95,"察",sz=130,b=True,c=GOODALL,a=PP_ALIGN.CENTER)
tb(s,2.85,4.10,2.30,0.30,"chá (examine)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
panel3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel3.fill.solid(); panel3.fill.fore_color.rgb=WHITE
panel3.line.color.rgb=GOODALL; panel3.line.width=Pt(2.5)
head3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head3.fill.solid(); head3.fill.fore_color.rgb=GOODALL; head3.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写  How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"1️⃣ 「观」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  左边「又」, 右边「见」(看)",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"2️⃣ 「察」 — 14 笔 (复杂!)",sz=14,b=True,c=DARK)
tb(s,5.45,3.40,4.10,0.30,"  上面「宀」, 下面「祭」",sz=10,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格练 3 遍",sz=12,b=True,c=GOODALL)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「观察」! 我想 观察 ___ 。",
    "I can write 观察! I want to observe ___.")
n+=1; pn(s,n)
notes(s,"6-7 分钟 — 「察」 是复杂字 (14 笔):\n• K 学生: 只写 「观」 就行\n• G3+: 两个都写\n• 提示: 「观」右边是「见」 — 用眼睛看! 「察」上面像房顶, 在房子里仔细看\n• 让学生说想观察什么 (蚂蚁? 弟弟? 天气?)\n• 完成 → 「我会写」 贴纸, 写得最好的 → 队 +2 分")

# ============================================================
# 51. SESSION 3 DIVIDER
# ============================================================
s=div("Session 3  下午 2:00–2:50","🧪 实验课  Engineering Challenge  ·  50 min",ENG,"🛠️"); n+=1; pn(s,n)

# ============================================================
# === SESSION 3 · PHASE 1: HOOK (5 min) ===
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🌉 怎么 过 去?  How to Cross? · 纸桥挑战",LAB)
# Two books with gap visual (left side)
book1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(2.40),Inches(2.0),Inches(1.20))
book1.fill.solid(); book1.fill.fore_color.rgb=BROWN; book1.line.color.rgb=DARK; book1.line.width=Pt(3)
tb(s,0.5,2.85,2.0,0.5,"📚",sz=42,a=PP_ALIGN.CENTER)
book2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.30),Inches(2.40),Inches(2.0),Inches(1.20))
book2.fill.solid(); book2.fill.fore_color.rgb=BROWN; book2.line.color.rgb=DARK; book2.line.width=Pt(3)
tb(s,4.30,2.85,2.0,0.5,"📚",sz=42,a=PP_ALIGN.CENTER)
# Gap with question mark
tb(s,2.50,2.50,1.80,0.5,"⬅️ 10 cm ➡️",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,2.50,2.95,1.80,0.7,"❓",sz=44,b=True,c=LAB,a=PP_ALIGN.CENTER)
# Tiny student trying to cross
tb(s,2.50,3.65,1.80,0.4,"小朋友想过去!",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
# Right side: paper challenge prompt
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.55),Inches(2.10),Inches(3.10),Inches(2.10))
panel.fill.solid(); panel.fill.fore_color.rgb=ENG; panel.line.color.rgb=YELLOW; panel.line.width=Pt(3)
tb(s,6.55,2.20,3.10,0.5,"📄 1 张 纸",sz=20,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,6.55,2.70,3.10,0.4,"1 sheet of paper",sz=11,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,6.55,3.10,3.10,0.5,"➕",sz=24,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,6.55,3.50,3.10,0.5,"= 一座桥?",sz=20,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,6.55,3.95,3.10,0.30,"= a bridge?",sz=11,c=WHITE,a=PP_ALIGN.CENTER)
# Hook prompt
tb(s,0.4,1.10,9.2,0.45,"❓ 只有 1 张 纸 — 怎么 造 一座 结实 的 桥?",sz=20,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,0.4,1.55,9.2,0.30,"Just 1 sheet of paper — how to build a strong bridge?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tpr_strip(s,4.45,"想一想 — 折? 卷? 撕? 全班一起想!","Think — fold? roll? tear? Whole class brainstorm!")
n+=1; pn(s,n)
notes(s,"🔥 HOOK · 5 分钟:\n• 老师摆 2 本书 (10 cm gap) — 让学生看 「桥要架在哪里」\n• 拿出 1 张 A4 纸 — 「就用这一张, 怎么造桥?」\n• 戏剧效果: 把纸平放架在书上 — 「这样行吗?」 (学生说不行/塌)\n• 引出问题: 「怎么让一张纸变结实? 折? 卷? 撕?」\n• 全班 30 秒头脑风暴 — 收集想法\n• 转: 「这就是工程师的工作 — 用普通材料造结实的东西!」")

# ============================================================
# === SESSION 3 · PHASE 2: MINI-LESSON (10 min) ===
# Engineering Design Process + Materials
# ============================================================
# Design process slide (cycle of 4)
s=ns(); bg(s,CREAM); hb(s,"🔁 工程师 4 步循环  Engineering Design Loop",ENG)
score_badge(s)
tb(s,0.4,0.85,9.2,0.4,"工程师从来不一次成功 — 改了再试!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.30,"Engineers never get it right the first time — they iterate!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("👀","看","Observe","脏水什么样?","SCI",SCI),
    ("💭","想","Design","怎么放材料?","ENG",ENG),
    ("🔨","做","Build","动手做!","RED",LAB),
    ("🔄","改","Improve","什么没用?","GREEN",GREEN),
]
# Arrange in a horizontal flow with arrows
for i,(em,cn,en,q,_,cl) in enumerate(steps):
    x=0.4+i*2.30
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.10),Inches(2.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(3)
    # Number badge
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.78),Inches(1.85),Inches(0.55),Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb=cl; nb.line.fill.background()
    tb(s,x+0.78,1.93,0.55,0.4,str(i+1),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.50,2.0,0.6,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.20,2.0,0.45,cn,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.0,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.95,2.0,0.30,q,sz=11,c=DARK,a=PP_ALIGN.CENTER)
    if i<3:
        tb(s,x+2.10,2.85,0.30,0.4,"→",sz=22,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
# Loop-back arrow
tb(s,0.4,4.40,9.2,0.30,"🔁 改完 → 再做 → 再改! (失败也是数据!)",sz=14,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,4.70,9.2,0.25,"Improve → Build → Improve again! (Failure is data too!)",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"📚 MINI 1 · 4 分钟:\n• 4 步走一遍 — 强调第 4 步「改」最重要\n• 「爱迪生试了 1000 次 — 失败 999 次也没放弃!」\n• 全班齐声: 看! 想! 做! 改! (加手势 — 看 = 指眼睛, 想 = 指头, 做 = 拳头, 改 = 转圈)")

# Materials slide — bridge challenge
s=ns(); bg(s,CREAM); hb(s,"🛠️ 材料 + 规则  Materials + Rules",ENG)
tb(s,0.4,0.85,9.2,0.4,"每队拿到一套材料 — 看清楚!",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
mats=[
    ("📄","A4 纸","1 sheet"),
    ("✂️","剪刀","Scissors"),
    ("📏","胶带","Tape"),
    ("📚","2 本书","2 books"),
    ("🪙","硬币","Coins"),
]
for i,(em,cn,en) in enumerate(mats):
    x=0.4+i*1.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.40),Inches(1.70),Inches(2.10))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=ENG; sh.line.width=Pt(2)
    tb(s,x,1.50,1.70,0.85,em,sz=50,a=PP_ALIGN.CENTER)
    tb(s,x,2.55,1.70,0.40,cn,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x,2.95,1.70,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Rules
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.65),Inches(9.2),Inches(1.50))
sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=ENG; sh.line.width=Pt(2)
tb(s,0.5,3.75,9.0,0.32,"📋 比赛规则  Competition Rules",sz=14,b=True,c=ENG)
tb(s,0.5,4.08,9.0,0.28,"1️⃣ 桥要架在 2 本书中间 (10 cm) · Bridge must span 10 cm gap",sz=11,c=DARK)
tb(s,0.5,4.36,9.0,0.28,"2️⃣ 试 2 次 — 第 1 次后必须 改! · 2 trials, must IMPROVE between",sz=11,c=DARK)
tb(s,0.5,4.64,9.0,0.28,"3️⃣ 老师加硬币 — 桥撑住最多硬币的队赢! +5 分",sz=11,c=DARK)
tb(s,0.5,4.92,9.0,0.20,"Teacher adds coins — team whose bridge holds most coins wins!",sz=9,c=GRAY)
n+=1; pn(s,n)
notes(s,"📚 MINI 2 · 6 分钟:\n• 1 分钟: 展示材料 — 让 1 个学生上台触摸\n• 2 分钟: 每队拿到自己的一套 (4-5 人一组, 4-6 队)\n• 1 分钟: 老师 demo — 把 1 张纸平放在 2 本书上 → 「行吗? (塌!) 那怎么办?」\n• 1 分钟: 强调规则 — 试 2 次 + 必须改\n• 1 分钟: 各队选「首席工程师」 (rotation)")

# ============================================================
# === SESSION 3 · PHASE 3: ACTIVE PRACTICE (20 min) ===
# Build → Test → IMPROVE → Build → Test → Compare
# ============================================================
# Round 1: Build (paper bridge)
s=ns(); bg(s,CREAM); hb(s,"🔨 第 1 轮 · 建 + 试  Round 1 · Build & Test",ENG)
tb(s,0.4,1.25,9.2,0.4,"⏱️ 5 分钟造桥 + 2 分钟测试",sz=20,b=True,c=ENG,a=PP_ALIGN.CENTER)
# Two halves
left=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.85),Inches(4.55),Inches(2.85))
left.fill.solid(); left.fill.fore_color.rgb=WHITE; left.line.color.rgb=ENG; left.line.width=Pt(3)
tb(s,0.4,1.95,4.55,0.45,"① 看 + 想 + 做",sz=18,b=True,c=ENG,a=PP_ALIGN.CENTER)
tb(s,0.5,2.45,4.35,0.30,"⏱ 5 分钟造桥",sz=12,b=True,c=GRAY,a=PP_ALIGN.LEFT)
tb(s,0.5,2.75,4.35,0.30,"• 看 — 怎么让纸变结实?",sz=12,c=DARK)
tb(s,0.5,3.05,4.35,0.30,"• 想 — 折? 卷? 加层?",sz=12,c=DARK)
tb(s,0.5,3.35,4.35,0.30,"• 做 — 全队一起动手",sz=12,c=DARK)
tb(s,0.5,3.65,4.35,0.30,"• 桥要架在 2 书中间!",sz=12,b=True,c=LAB)
tb(s,0.5,4.05,4.35,0.30,"💬 「我们的桥 ___ 。」",sz=11,b=True,c=ENG)
right=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.85),Inches(4.55),Inches(2.85))
right.fill.solid(); right.fill.fore_color.rgb=WHITE; right.line.color.rgb=LAB; right.line.width=Pt(3)
tb(s,5.05,1.95,4.55,0.45,"② 测试! Add Coins!",sz=18,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,5.15,2.45,4.35,0.30,"⏱ 2 分钟测试",sz=12,b=True,c=GRAY,a=PP_ALIGN.LEFT)
tb(s,5.15,2.75,4.35,0.30,"• 老师慢慢 加硬币",sz=12,c=DARK)
tb(s,5.15,3.05,4.35,0.30,"• 数: 1, 2, 3...",sz=12,c=DARK)
tb(s,5.15,3.35,4.35,0.30,"• 桥塌了 → 记下数字",sz=12,c=DARK)
tb(s,5.15,3.65,4.35,0.30,"• 哪里弱?",sz=12,b=True,c=GREEN)
tb(s,5.15,4.05,4.35,0.30,"💬 「撑住了 ___ 个硬币!」",sz=11,b=True,c=LAB)
sentence_frame_bar(s,4.85,"我们 ___ 纸 — 桥撑住 ___ 个硬币 — 我看到 ___ 。","We ___ the paper — bridge held ___ coins — I saw ___.",accent=ENG)
n+=1; pn(s,n)
notes(s,"🎯 ROUND 1 · 7 分钟:\n• 5 分钟造桥 — 老师巡视, 不给答案!\n• 2 分钟 test — 各队挨个测试 (老师加硬币)\n• 鼓励观察 — 「桥哪里先弯? 哪里先塌?」\n• 失败也好 — 这是数据! 「这告诉我们: 这种折法不够结实」\n• 记下每队的成绩 (硬币数), 准备 Round 2")

# Round 2: IMPROVE (key slide!)
s=ns(); bg(s,CREAM); hb(s,"🔄 关键一步 · 改!  KEY STEP · IMPROVE!",PURPLE)
tb(s,0.4,1.25,9.2,0.45,"工程师 = 改的高手!",sz=22,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,0.4,1.75,9.2,0.30,"Engineers = experts at IMPROVING!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# Three thinking prompts (bridge-specific)
prompts=[
    ("🤔","哪里弱?","Where is weak?","中间塌了?"),
    ("💡","怎么改?","How to fix?","折更多? 加柱子?"),
    ("🎯","目标!","Target!","撑更多硬币!"),
]
for i,(em,cn,en,detail) in enumerate(prompts):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.20),Inches(2.95),Inches(2.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=PURPLE; sh.line.width=Pt(3)
    tb(s,x+0.05,2.30,2.85,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,2.85,0.45,cn,sz=18,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.50,2.85,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.85,2.85,0.40,detail,sz=12,c=DARK,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.55,"第 1 次 ___ 没用 — 我要改成 ___ 。","First time ___ didn't work — I'll change to ___.",accent=PURPLE)
n+=1; pn(s,n)
notes(s,"🎯 IMPROVE · 3 分钟 (KEY moment!):\n• 「桥塌了 = 数据! 你刚学到了 — 这种方法不够结实」\n• 各队 30 秒讨论 — 「我们要改什么?」\n• 1 个代表上台说 (轮流上台) — 「我们要把 ___ 改成 ___」\n• 强调: 改一处就够 — 不是全部重做\n• 给每队 1 张新 A4 纸 (或重用旧的 — 老师决定)")

# Round 2: Build + Test (bridge)
s=ns(); bg(s,CREAM); hb(s,"🚀 第 2 轮 · 改了再试!  Round 2 · Try Again!",ENG)
tb(s,0.4,1.25,9.2,0.4,"⏱️ 5 分钟改 + 测试  ·  撑硬币最多的赢 +5 分!",sz=16,b=True,c=ENG,a=PP_ALIGN.CENTER)
# Compare visualization (weak bridge → strong bridge)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.85),Inches(9.2),Inches(2.85))
sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=ENG; sh.line.width=Pt(3)
# Trial 1: weak bridge (flat paper, sagging)
b1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(2.10),Inches(2.5),Inches(2.30))
b1.fill.solid(); b1.fill.fore_color.rgb=WARM; b1.line.color.rgb=DARK; b1.line.width=Pt(2)
tb(s,1.0,2.30,2.5,1.5,"📄",sz=72,a=PP_ALIGN.CENTER)
tb(s,1.0,3.55,2.5,0.35,"💔 5 个硬币",sz=14,b=True,c=LAB,a=PP_ALIGN.CENTER)
tb(s,1.0,3.85,2.5,0.30,"第 1 次 · Trial 1",sz=11,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1.0,4.10,2.5,0.22,"flat — sags",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Arrow with IMPROVE label
tb(s,3.55,2.50,1.20,0.40,"改 → ",sz=20,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,3.55,2.90,1.20,0.40,"IMPROVE",sz=12,b=True,c=PURPLE,a=PP_ALIGN.CENTER)
tb(s,3.55,3.30,1.20,0.30,"折/卷/加柱",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Trial 2: strong bridge (folded/curved - better)
b2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(2.10),Inches(2.5),Inches(2.30))
b2.fill.solid(); b2.fill.fore_color.rgb=YELLOW; b2.line.color.rgb=DARK; b2.line.width=Pt(3)
tb(s,4.85,2.30,2.5,1.5,"🌉",sz=72,a=PP_ALIGN.CENTER)
tb(s,4.85,3.55,2.5,0.35,"💪 20 个硬币!",sz=14,b=True,c=GREEN,a=PP_ALIGN.CENTER)
tb(s,4.85,3.85,2.5,0.30,"第 2 次 · Trial 2",sz=11,c=DARK,a=PP_ALIGN.CENTER)
tb(s,4.85,4.10,2.5,0.22,"folded — strong!",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
# Champion box
b3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.50),Inches(2.10),Inches(2.0),Inches(2.30))
b3.fill.solid(); b3.fill.fore_color.rgb=PH_ACTIVE; b3.line.fill.background()
tb(s,7.50,2.30,2.0,0.7,"🏆",sz=44,a=PP_ALIGN.CENTER)
tb(s,7.50,3.10,2.0,0.5,"冠军",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,7.50,3.55,2.0,0.30,"Champion",sz=11,c=YELLOW,a=PP_ALIGN.CENTER)
tb(s,7.50,3.95,2.0,0.30,"+5 分!",sz=14,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.85,"这次桥撑住 ___ 个硬币 — 因为我 ___ 。","This time bridge held ___ coins — because I ___.",accent=ENG)
n+=1; pn(s,n)
notes(s,"🎯 ROUND 2 · 8 分钟:\n• 5 分钟改 + 重新造桥\n• 3 分钟 test — 各队挨个测试, 加硬币\n• 各队的 2 次成绩对比 — Trial 1 vs Trial 2\n• 「谁撑得最多?」 — 冠军 +5 分\n• 颁奖 (3 个奖): 「最强桥」「最有创意」「最团结合作」 — 每队都有奖!\n• 重要: 不评失败 — 评进步! 「这队第 1 次 5 个, 第 2 次 15 个 → 进步最大也 +2 分!」")

# ============================================================
# === SESSION 3 · PHASE 4: APPLY (10 min) ===
# Reflection booklet + real-world tie-in
# ============================================================
# Reflection booklet page (5 min)
s=ns(); bg(s,CREAM); hb(s,"📔 我的工程笔记  My Engineer Notebook",GREEN)
group_label(s)
score_badge(s)
tb(s,0.4,1.25,9.2,0.4,"在小本子上写 / 画 — 3 行就够!",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.65,9.2,0.30,"Write or draw in your notebook — 3 lines is enough!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
lines=[
    ("👀","我看到 ___","I saw ___",SCI),
    ("🔧","我试了 ___","I tried ___",ENG),
    ("💡","我学到 ___","I learned ___",GREEN),
]
for i,(em,cn,en,cl) in enumerate(lines):
    y=2.10+i*0.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9.0),Inches(0.75))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2)
    tb(s,0.65,y+0.15,0.6,0.5,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.40,y+0.10,3.0,0.40,cn,sz=16,b=True,c=cl)
    tb(s,1.40,y+0.45,3.0,0.30,en,sz=10,c=GRAY)
    # Writing line
    line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(4.55),Inches(y+0.30),Inches(4.85),Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb=GRAY; line.line.fill.background()
n+=1; pn(s,n)
notes(s,"🌱 APPLY 1 · 5 分钟:\n• 老师发笔记本 / 已经在 booklet 里有这页\n• 个人写 5 分钟\n• K 学生 — 画就行 (老师可帮加字)\n• 老师巡视 — 找几个好的留到 share 用")

# Real-world bridges
s=ns(); bg(s,CREAM); hb(s,"🌍 真实世界 · 真的桥  Real-World Bridges",NAVY)
tb(s,0.4,1.25,9.2,0.4,"工程师造桥 — 让人安全过河、过山谷!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.65,9.2,0.30,"Engineers build bridges — helping people cross rivers, valleys safely!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
bridges=[
    ("🌉","金门大桥","Golden Gate Bridge","🇺🇸 美国 旧金山",WRIGHT),
    ("🚄","跨海大桥","Sea-crossing Bridge","🇨🇳 中国 港珠澳",ENG),
    ("🪵","古老石桥","Ancient stone bridges","几千年前 — 还在用!",NAVY),
    ("🏔️","山间吊桥","Mountain rope bridge","让人安全 cross",GREEN),
]
for i,(em,cn,en,where,cl) in enumerate(bridges):
    col=i%2; row=i//2
    x=0.4+col*4.7; y=2.10+row*1.30
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(1.20))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2)
    tb(s,x+0.10,y+0.10,0.95,1.0,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+1.15,y+0.15,3.3,0.40,cn,sz=15,b=True,c=cl)
    tb(s,x+1.15,y+0.55,3.3,0.30,en,sz=10,c=GRAY)
    tb(s,x+1.15,y+0.85,3.3,0.28,f"📍 {where}",sz=11,c=DARK)
n+=1; pn(s,n)
notes(s,"🌱 APPLY 2 · 5 分钟:\n• 介绍 4 座真桥 — 老师可以提前找图片\n• 重点: 金门大桥 — 「100 年前工程师就造出来了, 现在还在用!」\n• 港珠澳大桥 — 「世界最长跨海大桥」\n• 古老石桥 — 「几千年前的工程师 也用 折叠 / 拱形 让石头变结实!」\n• 「你今天做的实验 — 跟真工程师一样!」\n• 让 1-2 个学生说: 「我以后要造 ___」 (把今天和未来连起来)")

# ============================================================
# === SESSION 3 · PHASE 5: SHARE & CLOSE (5 min) ===
# Group show & tell + champion + badge
# ============================================================
s=ns(); bg(s,CREAM); hb(s,"🏆 颁奖 + Day 2 徽章!  Awards + Day 2 Badge!",PH_CLOSE)
score_badge(s)
# Champion announcement
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(1.45))
sh.fill.solid(); sh.fill.fore_color.rgb=PH_CLOSE; sh.line.color.rgb=DARK; sh.line.width=Pt(2.5)
tb(s,0.4,1.00,9.2,0.45,"🏆 今天总冠军队  Day Champions",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.50,9.2,0.45,"___ 队 — 全天最高分!",sz=24,b=True,c=ENG,a=PP_ALIGN.CENTER)
tb(s,0.4,2.00,9.2,0.30,"___ team — highest score across all 3 sessions!",sz=11,c=DARK,a=PP_ALIGN.CENTER)
# Badge
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.55),Inches(3),Inches(2.6))
sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=YELLOW; sh.line.width=Pt(6)
tf=tb(s,3.6,2.75,2.8,2.4,"DAY 2",sz=18,b=True,c=YELLOW,a=PP_ALIGN.CENTER)
ap(tf,"🔬👷",sz=42,a=PP_ALIGN.CENTER)
ap(tf,"小小科学家 & 工程师",sz=14,b=True,c=SCI,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=12,b=True,c=OK,a=PP_ALIGN.CENTER)
ap(tf,"💡🐒✈️🚦",sz=18,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"🎤 CLOSE · 5 分钟:\n• 1 分钟: 公布全天总分 — 冠军队上台 (轮流: 队长上)\n• 1 分钟: 4 队代表分享 — 「我们今天最喜欢 ___」 (30 秒/队, 轮流上台)\n• 2 分钟: 给每个学生发徽章 (or 贴纸) — 「你也是 problem solver!」\n• 1 分钟: 预告 Day 3 — 「明天 — 小小企业家! 乔布斯! 苹果!」\n• 全班合唱: 「科学家! 为什么? 工程师! 怎么办? 我也是 problem solver!」")

# ============================================================
out=os.path.join(os.path.dirname(__file__),"day2_scientists.pptx")
prs.save(out); print(f"Saved {out}  ({n} slides)")
