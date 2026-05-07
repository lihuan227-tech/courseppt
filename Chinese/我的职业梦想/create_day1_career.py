#!/usr/bin/env python3
"""
我的职业梦想 — Day 1: 认识职业世界 (Discover the World of Careers)
Same structure as 野外生存/小小艺术家 PBL, with a career-themed palette.
Show-don't-tell: students guess, act, choose, match, solve problems.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Career Adventure (navy + warm gold + helping green) ---
NAVY   = RGBColor(0x1E,0x3A,0x5F)   # primary — professional, trustworthy
GOLD   = RGBColor(0xF5,0xA6,0x23)   # secondary — energy, achievement
HELP   = RGBColor(0x2E,0x7D,0x32)   # accent — green for "helping"
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
BROWN  = RGBColor(0x6B,0x44,0x23)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
ALERT  = RGBColor(0xD0,0x4A,0x3C)
OK     = RGBColor(0x38,0x8E,0x3C)

# Job-specific colors
DOC    = RGBColor(0xE5,0x3E,0x3E)   # 医生 red
TEACH  = RGBColor(0x43,0xA0,0x47)   # 老师 green
POLICE = RGBColor(0x19,0x76,0xD2)   # 警察 blue
CHEF   = RGBColor(0xFB,0x8C,0x00)   # 厨师 orange
ENG    = RGBColor(0xFB,0xC0,0x2D)   # 工程师 yellow
ENV    = RGBColor(0x00,0x69,0x6B)   # 环境工程师 deep teal
WILD   = RGBColor(0x6D,0x4C,0x41)   # 野生动物保护员 brown
CITY   = RGBColor(0x7B,0x1F,0xA2)   # 城市规划师 purple

# === Helpers (same conventions as Day 2 camp / Day 1 art) ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])

def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt
    r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'
    return tf

def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt
    r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name='KaiTi'

def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)

def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=IMGBG; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)

def hb(s,txt,c=NAVY,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)

def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)

def notes(s,text):
    nf=s.notes_slide.notes_text_frame
    lines=text.split("\n")
    nf.text=lines[0]
    for line in lines[1:]:
        p=nf.add_paragraph(); p.text=line

def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)

def div(title,sub,color,emoji=""):
    s=ns(); bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    return s

def sentence_frame_bar(s,t,frame_cn,frame_en):
    # Clamp so bar (height 0.65) never overflows the 5.625" slide bottom
    if t > 4.95: t = 4.95
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid(); sf.fill.fore_color.rgb=WARM
    sf.line.color.rgb=GOLD; sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.7,0.4,"💬 我来说",sz=14,b=True,c=GOLD)
    tb(s,2.0,t+0.07,7.6,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.0,t+0.32,7.6,0.3,frame_en,sz=10,c=GRAY,a=None)

# === Specialized helpers ===

def mission_card(s,l,t,w,h,num,task_cn,task_en,emoji,color):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=color; sh.line.width=Pt(2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l+0.1),Inches(t+0.08),Inches(0.55),Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb=color; badge.line.fill.background()
    tb(s,l+0.1,t+0.18,0.55,0.4,str(num),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,l+0.7,t+0.15,w-0.8,0.4,task_en,sz=10,c=GRAY)
    tb(s,l+0.05,t+0.85,w-0.1,0.7,emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+1.55,w-0.1,0.4,task_cn,sz=18,b=True,c=color,a=PP_ALIGN.CENTER)

def guess_clue_slide(emoji,job_label,color,clues,frame_cn,frame_en):
    """Clue slide — students hear/see clues, guess the job (no answer shown)."""
    s=ns(); bg(s,CREAM)
    hb(s,f"{emoji} 猜猜我是谁？  Guess Who I Am!",color)
    tb(s,0.4,0.85,9.2,0.32,"听听我的提示 — 我是谁？",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.18,9.2,0.26,"Listen to my clues — who am I?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # LEFT — picture placeholder (mystery)
    img_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(4.30),Inches(2.85))
    img_box.fill.solid(); img_box.fill.fore_color.rgb=IMGBG
    img_box.line.color.rgb=color; img_box.line.width=Pt(2)
    tb(s,0.4,2.55,4.30,0.7,"❓",sz=70,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.40,4.30,0.40,"我是谁？",sz=18,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.80,4.30,0.30,"Who am I?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # RIGHT — clues
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.55),Inches(4.85),Inches(2.85))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
    panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.55),Inches(4.85),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.0,1.62,4.6,0.4,"🔎 提示 Clues",sz=14,b=True,c=WHITE)
    y=2.20
    for i,(c_cn,c_en) in enumerate(clues):
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(5.0),Inches(y+0.05),Inches(0.4),Inches(0.4))
        badge.fill.solid(); badge.fill.fore_color.rgb=color; badge.line.fill.background()
        tb(s,5.0,y+0.10,0.4,0.3,str(i+1),sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,5.5,y,4.1,0.4,c_cn,sz=14,b=True,c=DARK)
        tb(s,5.5,y+0.32,4.1,0.30,c_en,sz=9,c=GRAY)
        y+=0.65
    sentence_frame_bar(s,4.55,frame_cn,frame_en)
    return s

def reveal_job_slide(emoji,job_cn,job_en,color,what_cn,what_en,where_cn,where_en,video_url=""):
    """Reveal slide — show the job icon, name, what they do, where they work + video."""
    s=ns(); bg(s,CREAM)
    hb(s,f"{emoji} 答案揭晓!  Answer Revealed!",color)
    # Big emoji + job name (slightly shorter to make room for video bar)
    bigbox=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.30),Inches(3.05))
    bigbox.fill.solid(); bigbox.fill.fore_color.rgb=WHITE
    bigbox.line.color.rgb=color; bigbox.line.width=Pt(3)
    tb(s,0.4,1.20,4.30,1.40,emoji,sz=120,a=PP_ALIGN.CENTER)
    tb(s,0.4,2.65,4.30,0.55,job_cn,sz=28,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.20,4.30,0.40,job_en,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,1.4,3.65,2.3,0.30,"我就是 ___ !",GOLD,sz=11)
    # RIGHT — facts (also shorter)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(3.05))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
    panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    tb(s,5.0,1.10,4.6,0.40,"🛠️ 做什么? What I do",sz=15,b=True,c=color)
    tb(s,5.0,1.55,4.6,0.40,what_cn,sz=14,b=True,c=DARK)
    tb(s,5.0,1.95,4.6,0.30,what_en,sz=10,c=GRAY)
    tb(s,5.0,2.40,4.6,0.40,"📍 在哪里? Where",sz=15,b=True,c=color)
    tb(s,5.0,2.80,4.6,0.40,where_cn,sz=14,b=True,c=DARK)
    tb(s,5.0,3.20,4.6,0.30,where_en,sz=10,c=GRAY)
    # NEW: video bar full-width
    vid=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.10),Inches(9.30),Inches(0.55))
    vid.fill.solid(); vid.fill.fore_color.rgb=WHITE; vid.line.color.rgb=color; vid.line.width=Pt(1.5)
    tb(s,0.5,4.18,1.0,0.4,"📺 视频",sz=13,b=True,c=color)
    tb(s,1.5,4.18,8.1,0.30,video_url if video_url else "(老师在这里贴 1 分钟职业小视频)",sz=10,c=DARK)
    tb(s,1.5,4.45,8.1,0.20,"Watch ~1 min so kids see the job for real",sz=8,c=GRAY)
    sentence_frame_bar(s,4.85,
        f"我是 {job_cn} 。我帮助 ___ 。",
        f"I am a {job_en}. I help ___.")
    return s

def try_it_slide(emoji,job_cn,job_en,color,activity_cn,activity_en,
                 step1,step2,step3,say_cn,say_en):
    """1-2 minute hands-on activity right after a job reveal.
    Students literally do something the job does."""
    s=ns(); bg(s,CREAM)
    hb(s,f"🎬 试一试 · 当 {job_cn}!  Try It Like a {job_en}!",color)
    # Activity title banner
    title=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.20),Inches(0.7))
    title.fill.solid(); title.fill.fore_color.rgb=color; title.line.fill.background()
    tb(s,0.5,1.00,1.0,0.6,emoji,sz=32,a=PP_ALIGN.CENTER)
    tb(s,1.6,1.05,8.0,0.40,activity_cn,sz=20,b=True,c=WHITE)
    tb(s,1.6,1.42,8.0,0.28,activity_en,sz=11,c=WHITE)
    # 3 steps
    for i,step in enumerate([step1,step2,step3]):
        x=0.4+i*3.10
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(2.95),Inches(2.55))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=color; sh.line.width=Pt(2.5)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+1.225),Inches(1.95),Inches(0.50),Inches(0.50))
        badge.fill.solid(); badge.fill.fore_color.rgb=color; badge.line.fill.background()
        tb(s,x+1.225,2.04,0.50,0.35,str(i+1),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,x+0.10,2.55,2.75,0.40,step[0],sz=44,a=PP_ALIGN.CENTER)
        tb(s,x+0.10,3.10,2.75,0.45,step[1],sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.10,3.55,2.75,0.70,step[2],sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    # Time pill
    pill(s,4.0,4.55,2.0,0.30,"⏱️ 1-2 分钟  1-2 min",GOLD,sz=11)
    # Sentence frame
    sentence_frame_bar(s,4.95,say_cn,say_en)
    return s

def mystery_job_q_slide(emoji,job_label,color,picture_label,q1_cn,q1_en,q2_cn,q2_en,video_url=""):
    """Mystery job — picture + 2 simple A-or-B questions, NO label revealed yet."""
    s=ns(); bg(s,CREAM)
    hb(s,f"{emoji} 神秘职业  Mystery Job!",color)
    tb(s,0.4,0.85,9.2,0.32,"看看图 — 他们在做什么？",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.18,9.2,0.26,"Look — what are they doing?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # LEFT — bigger picture placeholder + video link below
    img_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.55),Inches(4.30),Inches(2.40))
    img_box.fill.solid(); img_box.fill.fore_color.rgb=IMGBG
    img_box.line.color.rgb=color; img_box.line.width=Pt(2.5)
    tb(s,0.4,2.05,4.30,0.55,emoji,sz=44,a=PP_ALIGN.CENTER)
    tb(s,0.4,2.65,4.30,0.35,picture_label,sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.05,4.30,0.30,"📷 在这里贴真实照片",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.40,4.30,0.30,"Paste a real photo here",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
    # Video link box under picture
    vid_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.05),Inches(4.30),Inches(0.65))
    vid_box.fill.solid(); vid_box.fill.fore_color.rgb=WHITE
    vid_box.line.color.rgb=color; vid_box.line.width=Pt(1.5)
    tb(s,0.5,4.10,1.0,0.4,"📺 视频",sz=12,b=True,c=color)
    tb(s,1.5,4.10,3.15,0.30,video_url if video_url else "(老师在这里贴视频链接)",sz=9,c=DARK)
    tb(s,1.5,4.40,3.15,0.25,"Watch ~1 min before discussion",sz=8,c=GRAY)
    # RIGHT — 2 questions
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.55),Inches(4.85),Inches(2.95))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
    panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.55),Inches(4.85),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.0,1.62,4.6,0.4,"🤔 想一想 Think First",sz=14,b=True,c=WHITE)
    badge1=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(5.0),Inches(2.30),Inches(0.45),Inches(0.45))
    badge1.fill.solid(); badge1.fill.fore_color.rgb=color; badge1.line.fill.background()
    tb(s,5.0,2.36,0.45,0.3,"1",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,5.55,2.25,4.05,0.40,q1_cn,sz=14,b=True,c=DARK)
    tb(s,5.55,2.62,4.05,0.30,q1_en,sz=9,c=GRAY)
    badge2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(5.0),Inches(3.20),Inches(0.45),Inches(0.45))
    badge2.fill.solid(); badge2.fill.fore_color.rgb=color; badge2.line.fill.background()
    tb(s,5.0,3.26,0.45,0.3,"2",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,5.55,3.15,4.05,0.40,q2_cn,sz=14,b=True,c=DARK)
    tb(s,5.55,3.52,4.05,0.30,q2_en,sz=9,c=GRAY)
    sentence_frame_bar(s,4.65,
        "我觉得他们 ___ 。",
        "I think they are ___.")
    return s

def mystery_job_label_slide(emoji,full_job_cn,full_job_en,kid_label_cn,kid_label_en,color,what_does):
    """Reveal the kid-friendly label for a mystery job."""
    s=ns(); bg(s,CREAM)
    hb(s,f"{emoji} 答案揭晓!  Answer!",color)
    # Big card
    big=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(0.95),Inches(8.8),Inches(2.7))
    big.fill.solid(); big.fill.fore_color.rgb=color; big.line.fill.background()
    tb(s,0.6,1.10,8.8,0.85,emoji,sz=70,a=PP_ALIGN.CENTER)
    tb(s,0.6,2.05,8.8,0.55,full_job_cn,sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.6,2.65,8.8,0.40,full_job_en,sz=14,c=WARM,a=PP_ALIGN.CENTER)
    # Kid-friendly label
    pill(s,2.5,3.10,5.0,0.50,kid_label_cn,GOLD,sz=18)
    tb(s,0.6,3.65,8.8,0.30,kid_label_en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    # What they do
    tb(s,0.6,4.05,8.8,0.40,what_does,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    sentence_frame_bar(s,4.65,
        f"{full_job_cn} 是 ___ 的人。",
        f"A {full_job_en} is a person who ___.")
    return s

def scenario_q_slide(em,scene_cn,scene_en,color,help_options):
    """Problem scenario — students decide who can help (3 options shown, no answer)."""
    s=ns(); bg(s,CREAM)
    hb(s,"🆘 谁来帮忙？  Who Can Help?",color)
    # scenario banner
    big=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(1.40))
    big.fill.solid(); big.fill.fore_color.rgb=color; big.line.fill.background()
    tb(s,0.4,1.05,9.2,0.50,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.55,9.2,0.45,scene_cn,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.4,2.00,9.2,0.30,scene_en,sz=11,c=WARM,a=PP_ALIGN.CENTER)
    # 3 choice cards (same 3 mystery jobs)
    for i,(j_em,j_cn,j_en,j_col) in enumerate(help_options):
        x=0.4+i*3.15
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.55),Inches(3.0),Inches(1.85))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
        sh.line.color.rgb=j_col; sh.line.width=Pt(2.5)
        tb(s,x+0.05,2.65,2.9,0.65,j_em,sz=38,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.38,2.9,0.40,j_cn,sz=15,b=True,c=j_col,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.80,2.9,0.30,j_en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
        pill(s,x+0.85,4.10,1.30,0.25,f"选 {chr(65+i)}",j_col,sz=10)
    sentence_frame_bar(s,4.55,
        "我选 ___ 来帮忙。",
        "I choose ___ to help.")
    return s

def scenario_a_slide(em,scene_cn,scene_en,answer_emoji,answer_cn,answer_en,color,why_cn,why_en):
    """Scenario reveal — show which job helps + 1-line why."""
    s=ns(); bg(s,CREAM)
    hb(s,"💡 答案揭晓!  Who Helps!",color)
    # scenario reminder
    big=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(9.2),Inches(0.85))
    big.fill.solid(); big.fill.fore_color.rgb=color; big.line.fill.background()
    tb(s,0.4,1.00,9.2,0.40,em+" "+scene_cn,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.4,1.42,9.2,0.30,scene_en,sz=10,c=WARM,a=PP_ALIGN.CENTER)
    # answer card (big)
    ans=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.0),Inches(2.05),Inches(6.0),Inches(2.30))
    ans.fill.solid(); ans.fill.fore_color.rgb=WHITE
    ans.line.color.rgb=color; ans.line.width=Pt(3)
    tb(s,2.0,2.15,6.0,0.85,answer_emoji,sz=64,a=PP_ALIGN.CENTER)
    tb(s,2.0,3.10,6.0,0.50,answer_cn,sz=24,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,2.0,3.65,6.0,0.35,answer_en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,3.5,4.00,3.0,0.30,"✅ 来帮忙!",OK,sz=11)
    # Why
    tb(s,0.4,4.45,9.2,0.30,f"💡 因为 {why_cn}",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.78,9.2,0.30,f"Because {why_en}",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    return s

# ============================================================
# === BUILD SLIDES ==========================================
# ============================================================
n=0

# 1. COVER
s=ns(); bg(s,NAVY)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.fill.background()
tb(s,1,0.4,8,0.5,"DAY 1",sz=18,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"💼 我的职业梦想",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"My Career Dream",sz=22,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"🌍 认识职业世界  Discover the World of Careers",sz=22,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"Guess · Act · Choose · Solve · Dream",sz=14,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"猜猜 · 演演 · 选选 · 解决 · 梦想",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"小小职业人 · Future Career Heroes",sz=14,b=True,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"开场 (1 分钟):\n• 「同学们! 长大后, 你想做什么?」\n• 今天我们去「职业世界」看看 — 不是听老师讲, 是自己猜、演、选!\n• 5 个步骤完成 → 拿「小小职业人」徽章!")

# 1.5  UNIT PREVIEW — 5-Day Career Journey
s=ns(); bg(s,CREAM); hb(s,"🗺️ 5 天的职业之旅  Our 5-Day Career Journey",NAVY)
tb(s,0.4,0.85,9.2,0.34,"5 天 · 5 个主题 · 每天认识一位/一群改变世界的人",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.28,"5 days · 5 themes · meet world-changers along the way",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
days_preview=[
    ("Day 1","认识职业世界","Discover Careers","🌍",NAVY,"今天 · 8 个职业"),
    ("Day 2","小小科学家","Little Scientists","🔬",ENV,"⭐ 爱因斯坦"),
    ("Day 3","小小企业家","Little Entrepreneurs","💡",GOLD,"⭐ 乔布斯"),
    ("Day 4","帮助别人","Helpers","❤️",DOC,"⭐ 医生 + 老师"),
    ("Day 5","AI 与未来","AI & the Future","🤖",CITY,"⭐ AI 公司 / 创始人"),
]
for i,(label,cn,en,em,cl,spotlight) in enumerate(days_preview):
    x=0.3+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.82),Inches(3.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.1),Inches(1.65),Inches(0.55),Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb=cl; badge.line.fill.background()
    tb(s,x+0.1,1.74,0.55,0.4,str(i+1),sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.7,1.7,1.1,0.3,label,sz=11,b=True,c=cl)
    tb(s,x+0.05,2.30,1.72,0.7,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.05,1.72,0.4,cn,sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,1.72,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    sep=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+0.25),Inches(3.85),Inches(1.32),Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb=cl; sep.line.fill.background()
    tb(s,x+0.05,4.00,1.72,0.85,spotlight,sz=11,b=True,c=cl,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,5.20,
    "我最想认识 ___ 。",
    "I most want to meet ___.")
n+=1; pn(s,n)
notes(s,"30 秒 — 只是预告:\n• 「这 5 天我们会认识各种各样的人 — 看看他们怎么改变世界!」\n• 不展开 — Day 1 内容下面继续。")

# 3. SESSION 1 DIVIDER
s=div("Session 1  上午 11:00–11:45","🎯 Career Adventure · 职业大冒险",NAVY,"🌟"); n+=1; pn(s,n)

# 4. STEP 1 — Warm-up: 你想当 ___ 吗?
s=ns(); bg(s,CREAM); hb(s,"👋 你想当...?  Do You Want to Be...?",GOLD)
tb(s,0.4,0.85,9.2,0.45,"举手 / 站起来 — 你想当吗？",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.30,"Raise your hand / stand up — do you want to be one?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
warm=[("🩺","医生","Doctor",DOC),
      ("📚","老师","Teacher",TEACH),
      ("👮","警察","Police",POLICE),
      ("👨‍🍳","厨师","Chef",CHEF),
      ("👷","工程师","Engineer",ENG)]
for i,(em,cn,en,c) in enumerate(warm):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(1.78),Inches(2.40))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.95,1.7,0.85,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.85,1.7,0.45,cn,sz=18,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.30,1.7,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.34,3.70,1.10,0.30,"举手 ✋",c,sz=10)
sentence_frame_bar(s,4.45,
    "我想当 ___ ! / 我不想当 ___ 。",
    "I want to be a ___! / I don't want to be ___.")
n+=1; pn(s,n)
notes(s,"5 分钟:\n• 老师指着每个职业: 「你想当医生吗?」让想当的孩子站起来或举手。\n• 不解释 — 直接玩!\n• 简单介绍: 职业 = 长大以后做的工作 (1 句话, 不要多)。\n• 让 1-2 个孩子说: 「我想当 ___ !」")

# 4.5  📖 PICTURE BOOK INTRO — "What Do People Do All Day?" by Richard Scarry
# Helper for picture book scene slides (3-level questions: see → job → problem)
def book_scene_slide(emoji,scene_cn,scene_en,color,
                     l1_q_cn,l1_q_en,l1_hint,
                     l2_q_cn,l2_q_en,l2_hint,
                     l3_q_cn,l3_q_en,l3_hint):
    s=ns(); bg(s,CREAM)
    hb(s,f"{emoji} 看故事 · {scene_cn}",color)
    # Big picture placeholder on left half
    img=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.30),Inches(4.10))
    img.fill.solid(); img.fill.fore_color.rgb=IMGBG; img.line.color.rgb=color; img.line.width=Pt(3)
    tb(s,0.4,2.10,4.30,1.2,emoji,sz=120,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.50,4.30,0.40,scene_cn,sz=18,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.92,4.30,0.30,scene_en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.45,4.30,0.30,"📷 在这里贴书中场景截图",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.72,4.30,0.25,"Paste book scene screenshot",sz=8,c=LGRAY,a=PP_ALIGN.CENTER)
    # Right column — 3 layered question cards (level 1, 2, 3)
    levels=[
        ("🟢","L1 你看到？  See",OK,l1_q_cn,l1_q_en,l1_hint),
        ("🟡","L2 是什么工作？  Job",GOLD,l2_q_cn,l2_q_en,l2_hint),
        ("🟠","L3 解决什么问题？  Solves",CHEF,l3_q_cn,l3_q_en,l3_hint),
    ]
    for i,(em,label,cl,q_cn,q_en,hint) in enumerate(levels):
        y=0.95+i*1.40
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(y),Inches(4.85),Inches(1.30))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
        head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(y),Inches(4.85),Inches(0.40))
        head.fill.solid(); head.fill.fore_color.rgb=cl; head.line.fill.background()
        tb(s,4.95,y+0.04,4.65,0.32,f"{em}  {label}",sz=12,b=True,c=WHITE)
        tb(s,4.95,y+0.45,4.65,0.35,q_cn,sz=14,b=True,c=DARK)
        tb(s,4.95,y+0.78,4.65,0.25,q_en,sz=9,c=GRAY)
        tb(s,4.95,y+1.02,4.65,0.25,f"💡 老师提示: {hint}",sz=9,c=cl)
    return s

# 4.5b — BOOK INTRO (绘本: 长大后你想做什么工作 — 小狗们的职业梦想)
s=ns(); bg(s,CREAM); hb(s,"📖 听一个故事  Story Time",GOLD)
tb(s,0.4,0.85,9.2,0.50,"📚 长大后你想做什么工作？",sz=22,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.4,1.40,9.2,0.30,"What Do You Want To Be When You Grow Up?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# Cover panel
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.85),Inches(4.30),Inches(3.20))
sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=GOLD; sh.line.width=Pt(3)
tb(s,0.4,2.10,4.30,1.0,"🐶",sz=110,a=PP_ALIGN.CENTER)
tb(s,0.4,3.15,4.30,0.45,"小狗们的职业梦想",sz=16,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.4,3.65,4.30,0.30,"The Puppies' Career Dreams",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,4.30,0.30,"📷 在这里贴书的封面",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.40,4.30,0.30,"Paste book cover here",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
# Right: video link + listening prompts
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.85),Inches(4.85),Inches(3.20))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=GOLD; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.85),Inches(4.85),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=GOLD; head.line.fill.background()
tb(s,5.00,1.93,4.65,0.4,"🎬 听故事  Listen",sz=14,b=True,c=WHITE)
tb(s,5.00,2.45,4.65,0.30,"📺 视频链接 / Video:",sz=11,b=True,c=GOLD)
tb(s,5.00,2.75,4.65,0.30,"youtube.com/watch?v=EVFPL_qXChU",sz=10,b=True,c=NAVY)
tb(s,5.00,3.05,4.65,0.30,"(老师: ~3-5 分钟, 看完一起聊)",sz=10,c=GRAY)
tb(s,5.00,3.50,4.65,0.28,"👂 听的时候, 注意 4 件事:",sz=11,b=True,c=GOLD)
tb(s,5.00,3.78,4.65,0.24,"1. 🐶 小狗们 想做 什么 工作?",sz=10,c=DARK)
tb(s,5.00,4.00,4.65,0.24,"2. ❤️ 它们 喜欢 什么? (兴趣)",sz=10,c=DARK)
tb(s,5.00,4.22,4.65,0.24,"3. 💪 它们 擅长 什么? (能力)",sz=10,c=DARK)
tb(s,5.00,4.44,4.65,0.24,"4. 🤝 它们 想 帮 谁? (帮助)",sz=10,c=DARK)
tb(s,5.00,4.70,4.65,0.28,"📌 听完, 一起 讨论!",sz=10,b=True,c=GOLD)
n+=1; pn(s,n)
notes(s,"3-5 分钟:\n• 主资源: https://www.youtube.com/watch?v=EVFPL_qXChU\n• 绘本:《长大后你想做什么工作》— 小狗们的职业梦想\n• 听之前先说: 「故事里的小狗都长大了 — 它们都想做什么工作? 一起听!」\n• 听完不要急着讲解 — 翻页让学生自己讨论\n• 关键: 这本书会带出 — 小狗们从 自己的兴趣 + 擅长 中找到职业 → 这就是今天的公式!")

# 4.5c — POST-STORY DISCUSSION (K-2 vs G3-5 differentiated questions with sample answers)
s=ns(); bg(s,CREAM); hb(s,"🤔 听完故事 · 一起讨论  Let's Talk",GOLD)
tb(s,0.4,0.80,9.2,0.32,"按年龄分 — K-2 / G3-5 不同问题",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.13,9.2,0.24,"Differentiated questions for K-2 vs G3-5",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# K-2 column (left)
k2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(1.45),Inches(4.55),Inches(3.50))
k2.fill.solid(); k2.fill.fore_color.rgb=WHITE; k2.line.color.rgb=DOC; k2.line.width=Pt(3)
k2h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(1.45),Inches(4.55),Inches(0.45))
k2h.fill.solid(); k2h.fill.fore_color.rgb=DOC; k2h.line.fill.background()
tb(s,0.50,1.50,4.40,0.36,"🐣 K-2  ·  小朋友 (简单)",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# K-2 Q1
tb(s,0.55,2.00,4.30,0.30,"❓ 1. 故事里的小狗, 长大都想做什么工作?",sz=11,b=True,c=DARK)
tb(s,0.55,2.30,4.30,0.22,"What jobs do the puppies want?",sz=8,c=GRAY)
tb(s,0.65,2.55,4.20,0.22,"💡 警察 · 医生 · 卖冰激凌 · 飞行员 ...",sz=10,b=True,c=DOC)
# K-2 Q2
tb(s,0.55,2.95,4.30,0.30,"❓ 2. 小狗们怎么决定做什么工作?",sz=11,b=True,c=DARK)
tb(s,0.55,3.25,4.30,0.22,"How did they decide their future jobs?",sz=8,c=GRAY)
tb(s,0.65,3.50,4.20,0.50,"💡 从生活经验、擅长的事、平日喜好里 联想出来的!",sz=10,b=True,c=DOC)
# K-2 sentence frame
sf1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(4.20),Inches(4.25),Inches(0.65))
sf1.fill.solid(); sf1.fill.fore_color.rgb=WARM; sf1.line.color.rgb=DOC; sf1.line.width=Pt(1.5)
tb(s,0.65,4.25,4.10,0.28,"💬 我喜欢的小狗想当 ___ 。",sz=11,b=True,c=DOC)
tb(s,0.65,4.55,4.10,0.22,"My favorite puppy wants to be a ___.",sz=8,c=GRAY)

# G3-5 column (right)
g3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.45),Inches(4.55),Inches(3.50))
g3.fill.solid(); g3.fill.fore_color.rgb=WHITE; g3.line.color.rgb=NAVY; g3.line.width=Pt(3)
g3h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.45),Inches(4.55),Inches(0.45))
g3h.fill.solid(); g3h.fill.fore_color.rgb=NAVY; g3h.line.fill.background()
tb(s,5.15,1.50,4.40,0.36,"🌟 G3-5  ·  大孩子 (深一点)",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# G3-5 Q1
tb(s,5.20,2.00,4.30,0.30,"❓ 1. 这本书想告诉我们什么道理?",sz=11,b=True,c=DARK)
tb(s,5.20,2.30,4.30,0.22,"What's the message of this book?",sz=8,c=GRAY)
tb(s,5.30,2.55,4.20,0.50,"💡 小事也能 启发 大梦想 — 兴趣 + 擅长 = 未来工作!",sz=10,b=True,c=NAVY)
# G3-5 Q2
tb(s,5.20,3.20,4.30,0.30,"❓ 2. 怎么把 兴趣 变成 工作?",sz=11,b=True,c=DARK)
tb(s,5.20,3.50,4.30,0.22,"How to turn interest into a career?",sz=8,c=GRAY)
tb(s,5.30,3.75,4.20,0.20,"💡 ① 兴趣 → 专业: 练到专业水准",sz=9,c=NAVY)
tb(s,5.30,3.95,4.20,0.20,"💡 ② 想 — 能不能 养活自己 (赚钱)?",sz=9,c=NAVY)
tb(s,5.30,4.15,4.20,0.20,"💡 ③ 多元 — 主业 + 副业 (斜杠)",sz=9,c=NAVY)
# G3-5 sentence frame
sf2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.20),Inches(4.45),Inches(4.25),Inches(0.45))
sf2.fill.solid(); sf2.fill.fore_color.rgb=WARM; sf2.line.color.rgb=NAVY; sf2.line.width=Pt(1.5)
tb(s,5.30,4.50,4.10,0.28,"💬 我的兴趣 ___ 可以变成 ___ 工作。",sz=10,b=True,c=NAVY)
tb(s,5.30,4.78,4.10,0.20,"My interest in ___ could become a ___ job.",sz=8,c=GRAY)
n+=1; pn(s,n)
notes(s,"6-8 分钟 (按年龄分组讨论):\n\n🐣 K-2 (低年级 — 具体, 简单):\n• Q1: 「故事里的小狗都想做什么工作?」让学生抢答 — 引导答案: 警察/医生/卖冰激凌/飞行员\n• Q2: 「他们是怎么想到的?」 — 答案: 从他们 喜欢的、擅长的、平时玩的东西 想出来的\n• 强调: 小狗也是从 自己的兴趣 找到 工作!\n\n🌟 G3-5 (大孩子 — 抽象, 深一点):\n• Q1 引出书的道理: 「小事可以启发大梦想 — 兴趣 + 擅长 = 工作」\n• Q2 引出实操: 怎么把兴趣变工作?\n  - ① 兴趣转专业 (要练到专业水准)\n  - ② 考虑能不能赚钱 (养活自己)\n  - ③ 不行就 斜杠 — 白天上班, 晚上做兴趣\n• 鼓励家长引导: 「什么事做起来 顺手? 什么事 特别喜欢?」\n\n• 不要混在一起讲 — 先 K-2 (3 分钟), 再 G3-5 (3 分钟)")

# 4.5c-2 — Key principle (moved here AFTER Discussion — students discover formula via discussion)
s=ns(); bg(s,GOLD)
tb(s,1,0.55,8,0.65,"🌟 一个大发现!  A Big Discovery!",sz=26,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.20,8,0.30,"Interest + Skill + Help = Career",sz=13,b=True,c=WARM,a=PP_ALIGN.CENTER)
# 4-element formula row: [兴趣] + [能力] + [帮助别人] = [职业]
fy=1.65    # formula y
bh=1.85    # box height
bw=1.85    # box width
opw=0.45   # operator width
fx=0.625   # formula left margin
formula=[
    ("Interest","兴趣","我喜欢的","I love ___","white",NAVY),
    ("Skill",   "能力","我擅长的","I'm good at ___","white",HELP),
    ("Help",    "帮助别人","我能帮的人","I help ___","white",DOC),
    ("Career",  "职业","我的职业!","my career!","navy",WHITE),
]
for i,(en,cn_lbl,cn_sub,en_sub,fill,accent) in enumerate(formula):
    bx = fx + i*(bw+opw)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(bx),Inches(fy),Inches(bw),Inches(bh))
    if fill=="navy":
        sh.fill.solid(); sh.fill.fore_color.rgb=NAVY; sh.line.color.rgb=GOLD; sh.line.width=Pt(4)
        text_c=WHITE; sub_c=WARM; en_c=GOLD
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=accent; sh.line.width=Pt(3)
        text_c=accent; sub_c=DARK; en_c=accent
    tb(s,bx+0.05,fy+0.10,bw-0.10,0.28,en,sz=11,b=True,c=en_c,a=PP_ALIGN.CENTER)
    tb(s,bx+0.05,fy+0.42,bw-0.10,0.65,cn_lbl,sz=24,b=True,c=text_c,a=PP_ALIGN.CENTER)
    tb(s,bx+0.05,fy+1.15,bw-0.10,0.30,cn_sub,sz=11,b=True,c=sub_c,a=PP_ALIGN.CENTER)
    tb(s,bx+0.05,fy+1.45,bw-0.10,0.25,en_sub,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
# Operators (+ + =)
for i,op in enumerate(["+","+","="]):
    ox = fx + (i+1)*bw + i*opw
    tb(s,ox,fy+0.55,opw,0.7,op,sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# Reflection prompt
tb(s,0.5,3.75,9,0.32,"💭 我的兴趣 + 我的能力 + 我想帮的人 = 什么职业?",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,4.10,9,0.26,"My interest + my skill + who I help = what career?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.5,4.80,9,0.34,"长大以后做什么? 可以从你的兴趣、能力和想帮助的人开始想。",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,5.20,9,0.24,"What do you want to do? Start with your interests, skills, and people you want to help.",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"2 分钟 — 关键时刻 (这个公式贯穿全天!):\n• 这是讨论后的总结 — 学生刚刚回答了 4 个问题, 现在揭示公式\n• 老师指着每个 box 解释:\n  - 兴趣 Interest: 你喜欢什么? (画画? 运动? 动物?)\n  - 能力 Skill: 你擅长什么? (跑得快? 算得快? 画得好?)\n  - 帮助 Help: 你想帮助谁? (病人? 小动物? 同学?)\n  - = 职业 Career: 三个加起来 = 你未来的工作!\n• 全班齐读 2 遍: 「兴趣 + 能力 + 帮助 = 职业!」\n• 这是 Day 1 的核心 — 后面所有讨论都回到这个公式!")

# 4.5d — INTEREST + SKILL + HELP → JOBS (4-col, formula-aligned, discussion-first)
s=ns(); bg(s,CREAM); hb(s,"🔗 兴趣 + 能力 + 帮助 → 职业  Interest + Skill + Help → Career",GOLD)
# Discussion prompt bar
prompt=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(0.65))
prompt.fill.solid(); prompt.fill.fore_color.rgb=NAVY; prompt.line.color.rgb=GOLD; prompt.line.width=Pt(2)
tb(s,0.55,0.90,9.0,0.28,"🤔 先讨论 — 一起猜! (don't read answers yet!)",sz=12,b=True,c=RGBColor(0xFF,0xC1,0x07),a=PP_ALIGN.CENTER)
tb(s,0.55,1.18,9.0,0.28,"❓ 兴趣 → ❓ 能力 → ❓ 帮助谁 → ❓ 可能的职业?",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# 4 column headers (4 col layout)
y0=1.62
hdr=[("🎨 兴趣 Interest",  GOLD, 0.40, 1.75),
     ("💪 能力 Skill",      HELP, 2.20, 1.95),
     ("🤝 帮助谁 Help",     CHEF, 4.20, 2.30),
     ("💼 可能的职业 Jobs", NAVY, 6.55, 3.05)]
for label,color,x,w in hdr:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y0),Inches(w),Inches(0.38))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background()
    tb(s,x,y0+0.04,w,0.30,label,sz=10,b=True,c=WHITE,a=PP_ALIGN.CENTER)
# 6 rows: (interest, skill, help, jobs)
mappings=[
    ("⚽","运动",     "跑得快 · 身体好 · 有耐心",     "让人健康 · 教孩子",        "运动员 · 教练 · 体育老师 · 体育记者"),
    ("🎨","画画",     "观察细 · 有创意 · 用色感",     "让生活美 · 给人快乐",       "画家 · 设计师 · 动画师 · 插画师"),
    ("🧱","搭积木",   "想象力 · 空间感 · 动手强",     "造楼 · 造桥 · 安全的家",    "工程师 · 建筑师 · 设计师"),
    ("🎮","玩游戏",   "反应快 · 解决问题 · 合作",     "让人开心 · 好玩的游戏",     "游戏设计师 · 程序员 · 测试员"),
    ("🐶","爱小动物", "耐心 · 细心 · 有爱心",         "帮生病动物 · 保护动物",     "兽医 · 动物保护员 · 训练师"),
    ("🍳","做饭",     "创造力 · 味觉好 · 动手强",     "做好吃的 · 家庭温暖",       "厨师 · 面包师 · 食品设计师"),
]
row_h=0.46
for i,(em,interest,skill,help_,jobs) in enumerate(mappings):
    y=2.05+i*row_h
    rh=row_h-0.05
    # Col 1: Interest
    c1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.40),Inches(y),Inches(1.75),Inches(rh))
    c1.fill.solid(); c1.fill.fore_color.rgb=WARM; c1.line.color.rgb=GOLD; c1.line.width=Pt(1.5)
    tb(s,0.45,y+0.05,0.50,0.30,em,sz=16,a=PP_ALIGN.CENTER)
    tb(s,1.00,y+0.08,1.10,0.30,interest,sz=12,b=True,c=DARK)
    # Col 2: Skill
    c2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.20),Inches(y),Inches(1.95),Inches(rh))
    c2.fill.solid(); c2.fill.fore_color.rgb=WHITE; c2.line.color.rgb=HELP; c2.line.width=Pt(1.5)
    tb(s,2.28,y+0.08,1.85,0.30,skill,sz=8,b=True,c=HELP)
    # Col 3: Help
    c3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.20),Inches(y),Inches(2.30),Inches(rh))
    c3.fill.solid(); c3.fill.fore_color.rgb=WHITE; c3.line.color.rgb=CHEF; c3.line.width=Pt(1.5)
    tb(s,4.28,y+0.08,2.20,0.30,help_,sz=8,b=True,c=CHEF)
    # Col 4: Jobs
    c4=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.55),Inches(y),Inches(3.05),Inches(rh))
    c4.fill.solid(); c4.fill.fore_color.rgb=WHITE; c4.line.color.rgb=NAVY; c4.line.width=Pt(1.5)
    tb(s,6.63,y+0.08,2.95,0.30,jobs,sz=8,b=True,c=NAVY)
n+=1; pn(s,n)
notes(s,"8-10 分钟 (DISCUSS FIRST — 不要急着读答案!):\n• 老师指着第一行 ⚽ 运动 — 一列一列问:\n  - Q1 兴趣: 「喜欢运动的人, 有什么 能力?」 (跑得快? 身体好?)\n  - Q2 帮助: 「他们能帮谁?」 (让人健康)\n  - Q3 职业: 「可以变什么工作?」 (运动员? 教练?)\n• 让 2-3 个学生猜每一列, 再揭示\n• 一行一行做 (~1.5 分钟/行)\n• 关键 takeaway: 一个兴趣 + 不同能力 → 很多不同的工作!")

# 4.5d-2 — TURN & TALK · 4 sequential rounds (Interest → Skill → Help → Career)
# Helper for one Turn & Talk round slide
def tt_round_slide(round_num, em, theme_cn, theme_en, q_cn, q_en, frame_cn, frame_en,
                   samples, color, bonus_lines=None):
    """One round of Turn & Talk — single question, big & focused."""
    s=ns(); bg(s,CREAM); hb(s,f"🗣️ Turn & Talk · Round {round_num} · {theme_cn}",color)
    # Big round badge + question card
    qcard=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.85),Inches(9.2),Inches(1.50))
    qcard.fill.solid(); qcard.fill.fore_color.rgb=color; qcard.line.fill.background()
    # Round badge (left)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.65),Inches(1.10),Inches(1.0),Inches(1.0))
    badge.fill.solid(); badge.fill.fore_color.rgb=WHITE; badge.line.color.rgb=color; badge.line.width=Pt(3)
    tb(s,0.65,1.30,1.0,0.6,str(round_num),sz=42,b=True,c=color,a=PP_ALIGN.CENTER)
    # Big emoji + question
    tb(s,1.85,1.00,7.5,0.85,em,sz=48,a=PP_ALIGN.LEFT)
    tb(s,2.85,1.05,6.6,0.55,q_cn,sz=26,b=True,c=WHITE)
    tb(s,2.85,1.65,6.6,0.30,q_en,sz=12,c=WARM)
    tb(s,2.85,1.95,6.6,0.30,theme_en,sz=10,c=WARM)
    # Sentence frame strip
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(2.50),Inches(9.2),Inches(0.65))
    sf.fill.solid(); sf.fill.fore_color.rgb=WARM; sf.line.color.rgb=color; sf.line.width=Pt(2.5)
    tb(s,0.55,2.55,2.0,0.30,"💬 句型 Frame",sz=12,b=True,c=color)
    tb(s,2.50,2.55,7.0,0.32,frame_cn,sz=18,b=True,c=DARK)
    tb(s,2.50,2.88,7.0,0.22,frame_en,sz=10,c=GRAY)
    # Teacher samples (4 examples in a row)
    tb(s,0.4,3.25,9.2,0.30,"💡 老师举例 · Teacher Examples (don't read all — pick 2-3)",sz=11,b=True,c=color,a=PP_ALIGN.CENTER)
    n_samples=len(samples)
    sw=(9.2-0.20*(n_samples-1))/n_samples
    for i,(s_em,s_cn) in enumerate(samples):
        x=0.4+i*(sw+0.20)
        sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.60),Inches(sw),Inches(0.85))
        sb.fill.solid(); sb.fill.fore_color.rgb=WHITE; sb.line.color.rgb=color; sb.line.width=Pt(2)
        tb(s,x,3.70,sw,0.40,s_em,sz=24,a=PP_ALIGN.CENTER)
        tb(s,x,4.10,sw,0.32,s_cn,sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # Bonus / extension lines at bottom
    if bonus_lines:
        bb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.60),Inches(9.2),Inches(0.55))
        bb.fill.solid(); bb.fill.fore_color.rgb=color; bb.line.fill.background()
        for i,line in enumerate(bonus_lines[:2]):
            tb(s,0.55,4.65+i*0.25,9.0,0.25,line,sz=11,b=True,c=WHITE)
    return s

# Round 1 — Interest (easiest, broadest entry)
s=tt_round_slide(1,"❤️","兴趣 Interest","Start with what you LOVE",
    "你喜欢什么?","What do you love?",
    "我喜欢 ____ 。","I love ____.",
    [("🎨","画画"),("🐶","动物"),("⚽","运动"),("📚","读书")],
    GOLD,
    bonus_lines=["💭 加难: 为什么? Why?  →  「因为很好玩」「因为我觉得很酷」"])
n+=1; pn(s,n)
notes(s,"⏱️ 1 分钟:\n• 一次只问 一个问题! (这是 Round 1 — 最简单)\n• 让学生和同桌轮流说: 「我喜欢 ___」\n• 老师举 2-3 个 例子 — 让 K 学生 也敢说\n• 高年级: 加 「为什么?」 — 因为好玩 / 因为 cool\n• 不评价 — 鼓励所有答案")

# Round 2 — Skill (with Growth Mindset)
s=tt_round_slide(2,"💪","能力 Skill","What you're good at",
    "你擅长什么?","What are you good at?",
    "我擅长 ____ 。","I'm good at ____.",
    [("🎨","画画"),("🏃","跑步"),("🧱","乐高"),("🧮","数学"),("📖","讲故事")],
    HELP,
    bonus_lines=["🌱 现在 不会, 也可以练习! Not good yet? You can practice!",
                 "💬 加: 我正在学习 ____ 。  I'm working on ____."])
n+=1; pn(s,n)
notes(s,"⏱️ 1.5 分钟 — 这是关键的一轮!\n• Round 2 比 Round 1 难一点 — 学生可能不知道 「擅长什么」\n• 给 choices: 画画 / 跑步 / 乐高 / 数学 / 讲故事 / 帮助别人\n• ❗ 重要: 告诉学生 「现在不会, 也可以练习!」 — Growth Mindset!\n• 如果学生说 「我不擅长任何东西」: 让他们说 「我正在学习 ___」\n• 这个 frame 特别 powerful — 让学生 own 自己的成长\n• 翻下一页: 深入 Growth Mindset")

# Round 2.5 — GROWTH MINDSET (能力 是 练 出来 的) — moved here to follow Round 2 Skill
s=ns(); bg(s,CREAM); hb(s,"💪 能力 是 练 出来 的!  Skills Are Built Through Practice",HELP)
tb(s,0.4,0.85,9.2,0.40,"现在 不会 — 没关系! 我们 都在 练习!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Not good at it yet? That's OK! We are all working on our skills.",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 3 reflection prompts (3-card row)
gm=[
    ("🤔","这个工作 需要 什么 能力?","What skills does this job need?",HELP),
    ("🌱","我 现在 还 不会 ___ ,","I'm not very good at ___ yet,",NAVY),
    ("💪","但我 正在 练习!","but I am working on this skill!",GOLD),
]
for i,(em,cn,en,c) in enumerate(gm):
    x=0.4+i*3.10
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.75),Inches(2.95),Inches(2.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,1.85,2.85,0.85,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,2.85,0.65,cn,sz=14,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,2.85,0.55,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# Sentence frames (whole class)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.35),Inches(9.2),Inches(1.05))
sf.fill.solid(); sf.fill.fore_color.rgb=HELP; sf.line.color.rgb=GOLD; sf.line.width=Pt(2.5)
tb(s,0.55,4.42,9.0,0.30,"💬 我 正在 努力 练习 ___ 。",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,4.74,9.0,0.24,"I am working on ___.",sz=10,c=WARM,a=PP_ALIGN.CENTER)
tb(s,0.55,5.02,9.0,0.30,"💬 我 现在 还 不太 会 ___ , 但是 我 正在 练习。",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — 成长型思维 (Growth Mindset) — 紧接着 Round 2 Skill:\n• 关键 message: 能力 不是 天生 的 — 是 练 出来 的!\n• 老师可以问: 「这个工作需要什么 能力?」 (例: 医生 — 耐心, 细心, 喜欢 帮助 人)\n• 让学生想: 「我现在 不擅长 什么? 我可以 练什么?」\n• 用句型说: 「我 现在 还 不太 会 ___, 但是 我 正在 练习」\n• 表扬: 「敢说 不会 — 才能 学会!」\n• 联系下面的 5 个职业 — 每个都需要 不同 的 能力 + 练习")

# Round 3 — Help
s=tt_round_slide(3,"🤝","帮助 Help","Who you want to help",
    "你想帮助谁?","Who do you want to help?",
    "我想帮助 ____ 。","I want to help ____.",
    [("🐾","小动物"),("🤒","生病的人"),("🧒","小朋友"),("🌍","地球"),("👨‍👩‍👧","爸爸妈妈")],
    DOC,
    bonus_lines=["✨ 让职业 不只是 赚钱 — 是 帮人!",
                 "Career is more than money — it's about helping people!"])
n+=1; pn(s,n)
notes(s,"⏱️ 1 分钟:\n• 「你想帮助谁?」 — 这一轮 让职业 有意义\n• 引导多元答案: 小动物 / 生病的人 / 小朋友 / 地球 / 爸爸妈妈\n• 重点: 不要 让职业 只关于 钱 — 强调 帮人")

# Round 4 — Career (connect everything + share)  · uses CHAINED example sentence
s=tt_round_slide(4,"💼","职业 Career","Putting it all together",
    "你长大想做什么?","What do you want to be?",
    "我喜欢___, 我擅长___, 我想帮___, 所以我想当___!",
    "I love ___, I'm good at ___, I want to help ___, so I want to be a ___!",
    [("🩺","医生"),("👨‍🏫","老师"),("👩‍🚀","宇航员"),("🎨","动画师"),("🦒","兽医")],
    NAVY,
    bonus_lines=["📝 例子 Example: 我喜欢画画, 我擅长画卡通, 我想帮小朋友开心, 所以我想当动画师!",
                 "🎤 找 2-3 个学生上台分享 (Whole-Class Share)"])
n+=1; pn(s,n)
notes(s,"⏱️ 2 分钟 (最后一轮 + 全班分享):\n• Round 4 把前面 3 轮 串起来 — 用 完整 chain sentence\n• 老师 念 example: 「我喜欢画画, 我擅长画卡通, 我想帮小朋友开心, 所以我想当动画师!」\n• 让 2-3 个学生上台 用同样的 chain:\n  - 「我喜欢动物, 我擅长安静观察, 我想帮 流浪猫, 所以我想当兽医!」\n  - 「我喜欢运动, 我擅长跑步, 我想帮 小朋友 健康, 所以我想当 体育老师!」\n• Bonus question: 「这个工作 会帮助谁?」 — 让所有职业 都关于帮人\n• 收尾金句: 「🌟 每个人都可以用 自己的兴趣 和 能力 帮助 世界」\n• 这个金句 是 整套课的 灵魂!")

# 4.5d-3 — INTEREST SURVEY (movement activity, energetic)
s=ns(); bg(s,CREAM); hb(s,"🎮 兴趣大调查  Interest Survey · Move with Me!",NAVY)
tb(s,0.4,0.85,9.2,0.40,"老师说 — 你做! 听到你喜欢的, 一起动起来!",sz=18,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Teacher says — YOU do! Hear your interest? Move!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 4 prompt cards (2x2 grid)
prompts=[
    ("⚽","谁喜欢运动?","Who loves sports?",       "站起来!","STAND UP!",  "→ 你可能是未来的运动员!","Future athlete!",   NAVY),
    ("🎨","谁喜欢画画?","Who loves drawing?",      "举手!",  "HAND UP!",   "→ 你可能是未来的设计师!","Future designer!",  GOLD),
    ("🐶","谁喜欢动物?","Who loves animals?",      "拍拍手!","CLAP CLAP!", "→ 你可能是未来的兽医!",  "Future vet!",       HELP),
    ("🤖","谁喜欢机器人?","Who loves robots?",     "跳一跳!","JUMP!",      "→ 你可能是未来的工程师!","Future engineer!",  CITY),
]
for i,(em,q_cn,q_en,act_cn,act_en,career_cn,career_en,color) in enumerate(prompts):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.75+row*1.75
    # Big card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.50),Inches(1.65))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=color; sh.line.width=Pt(3)
    # Big emoji on left
    tb(s,x+0.10,y+0.15,1.10,1.0,em,sz=44,a=PP_ALIGN.CENTER)
    # Question (top right)
    tb(s,x+1.30,y+0.10,3.10,0.32,q_cn,sz=14,b=True,c=color)
    tb(s,x+1.30,y+0.40,3.10,0.22,q_en,sz=9,c=GRAY)
    # Action prompt (highlighted bar)
    actrect=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+1.30),Inches(y+0.65),Inches(3.10),Inches(0.42))
    actrect.fill.solid(); actrect.fill.fore_color.rgb=color; actrect.line.fill.background()
    tb(s,x+1.30,y+0.71,3.10,0.30,act_cn,sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # Career hint at bottom of card
    tb(s,x+0.10,y+1.18,4.30,0.24,career_cn,sz=11,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,y+1.40,4.30,0.18,career_en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"3-4 分钟 (能量满满!):\n• 老师 喊一个 prompt, 学生立刻 做动作\n• 节奏: 每个 prompt = 5-10 秒, 等所有反应后再下一个\n• 看到学生反应后, 老师点出对应 career: 「太棒了! 你可能是未来的运动员!」\n• 4 轮做完, 学生坐下\n• 关键: 把 兴趣 + 工作 连起来 — 让学生感受 「我也可以!」")

# (Growth Mindset slide moved earlier — now appears right after Round 2 Skill)

# 7-16. Five GUESS clue + reveal pairs
# Clue 难度 progression: 兴趣 (抽象) → 擅长 (中等) → 具体 hint (近答案)
# Clue 1 不能 直接 说 职业 — 要让学生 听 完 3 个 才猜!
guesses=[
    ("🩺","Doctor",DOC,
     [("❤️ 我的兴趣是: 让 别人 不 难受。","❤️ My interest: making others feel better."),
      ("💪 我擅长: 找出 哪里 不 对劲。","💪 I'm good at: finding what's wrong."),
      ("👕 我穿 白 大褂, 用 听诊器 听 心跳。","👕 I wear a white coat, use a stethoscope.")],
     "医生","Doctor","看病、给药、做手术","check, medicine, surgery","医院 / 诊所","hospital / clinic",
     "搜: 'doctor for kids' (Nat Geo Kids / SciShow Kids)"),
    ("📚","Teacher",TEACH,
     [("❤️ 我的兴趣是: 看 别人 学会 新 东西。","❤️ My interest: seeing others learn something new."),
      ("💪 我擅长: 把 难的 变 简单。","💪 I'm good at: making hard things simple."),
      ("🏫 我每天 在 学校 跟 你 说话, 还会 改 作业。","🏫 I talk with you at school every day, and grade homework.")],
     "老师","Teacher","教课、改作业、讲故事","teach, grade, tell stories","学校 / 教室","school / classroom",
     "搜: 'a day in the life of a teacher for kids'"),
    ("👮","Police",POLICE,
     [("❤️ 我的兴趣是: 让 坏事 不 发生。","❤️ My interest: stopping bad things from happening."),
      ("💪 我擅长: 跑得 快, 听到 哨子 立刻 行动。","💪 I'm good at: running fast, acting on a whistle."),
      ("👮 我穿 制服, 开 红蓝灯 的 车。","👮 I wear a uniform, drive a red-blue-light car.")],
     "警察","Police Officer","抓坏人、保护大家","catch bad guys, protect people","街上 / 警察局","street / police station",
     "搜: 'police officer for kids' (Cocomelon / Sesame Street)"),
    ("👨‍🍳","Chef",CHEF,
     [("❤️ 我的兴趣是: 让 大家 笑着 吃 完。","❤️ My interest: making people smile while they eat."),
      ("💪 我擅长: 闻 味道, 配 颜色, 用 火。","💪 I'm good at: smelling flavors, mixing colors, using fire."),
      ("🍳 我戴 高高的 白 帽子, 在 厨房 里。","🍳 I wear a tall white hat, in the kitchen.")],
     "厨师","Chef","做菜、做点心","cook meals, make desserts","餐厅 / 厨房","restaurant / kitchen",
     "搜: 'kid chef cooking' or 'MasterChef Junior' (1-min clip)"),
    ("👷","Engineer",ENG,
     [("❤️ 我的兴趣是: 把 想 出来 的 东西 变 真。","❤️ My interest: turning ideas into real things."),
      ("💪 我擅长: 画 图, 量 尺寸, 解决 难题。","💪 I'm good at: drawing plans, measuring, solving puzzles."),
      ("🏗️ 我 画的图 后来 变 成 大楼、大桥。","🏗️ My drawings later become buildings & bridges.")],
     "工程师","Engineer","设计、建造、修理","design, build, repair","公司 / 工地","company / construction site",
     "搜: 'what does an engineer do for kids' (Crash Course Kids)"),
]
try_it_data={
    # All individual Try-It slides removed per user request:
    #   - 医生/老师 deleted (slides 13/16)
    #   - 警察/厨师 deleted earlier
    #   - 工程师 moved to Session 3 as a team paper-bridge competition
}
for em,j_en,col,clues,job_cn,job_en,what_cn,what_en,where_cn,where_en,video in guesses:
    s=guess_clue_slide(em,j_en,col,clues,
        f"我猜是 ___ !",
        f"I guess it's a ___!")
    n+=1; pn(s,n)
    notes(s,f"猜谜 (1-2 分钟):\n• 老师慢慢读 2 个提示, 让学生先想。\n• 让 2-3 个学生举手猜: 「我猜是 ___ !」\n• 不要直接给答案 — 让多个孩子先猜。\n• 然后翻页揭晓!")

    s=reveal_job_slide(em,job_cn,j_en,col,what_cn,what_en,where_cn,where_en,video_url=video)
    n+=1; pn(s,n)
    notes(s,f"揭晓 {job_cn}!\n• 让全班一起喊: 「{job_cn}!」\n• 看 1 分钟职业小视频 (老师提前在 YouTube 找好)\n• 复习: 做什么? 在哪里?\n• 用句型: 「我是 {job_cn}, 我帮助 ___ 。」")

    # NEW: Try-it micro-activity right after reveal
    if job_cn in try_it_data:
        a_cn,a_en,st1,st2,st3,say_cn,say_en=try_it_data[job_cn]
        s=try_it_slide(em,job_cn,j_en,col,a_cn,a_en,st1,st2,st3,say_cn,say_en)
        n+=1; pn(s,n)
        notes(s,f"试一试 (~2 分钟):\n• 让全班真的做一遍 — 不只看, 要做!\n• 给积极参与的孩子一个 ✓ 在白板上。\n• 简单收尾: 「现在你也是小小 {job_cn}!」")

# 16.5. COOL JOBS SHOWCASE — wide variety of "wow" jobs
s=ns(); bg(s,CREAM); hb(s,"🌟 还有这些超酷的职业!  More Cool Careers!",GOLD)
tb(s,0.4,0.85,9.2,0.34,"世界上有几千种职业 — 看看这 9 个超酷的!",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.20,9.2,0.26,"There are thousands of jobs — here are 9 cool ones!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
cool_jobs=[
    ("🚀","宇航员","Astronaut","在太空漂!",NAVY),
    ("✈️","飞行员","Pilot","开飞机看云!",POLICE),
    ("🦒","兽医","Vet","给小动物看病!",WILD),
    ("🎮","游戏设计师","Game Designer","你玩的游戏 = 他做的!",CITY),
    ("🌊","海洋生物学家","Marine Biologist","和海豚做朋友!",ENV),
    ("🚒","消防员","Firefighter","坐红色大车救人!",DOC),
    ("🎬","动画师","Animator","你看的卡通 = 他画的!",CHEF),
    ("🍫","巧克力师","Chocolatier","每天发明新糖果!",GOLD),
    ("🎩","魔术师","Magician","让东西「消失」!",HELP),
]
for i,(em,cn,en,wow,cl) in enumerate(cool_jobs):
    col=i%3; row=i//3
    # 3×3 grid — tighter cards to fit 9
    x=0.30+col*3.20; y=1.45+row*1.20
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.05))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,x+0.10,y+0.07,0.7,0.55,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.08,2.10,0.32,cn,sz=14,b=True,c=cl)
    tb(s,x+0.85,y+0.36,2.10,0.24,en,sz=9,c=GRAY)
    tb(s,x+0.10,y+0.65,2.85,0.35,f"💫 {wow}",sz=10,c=DARK)
n+=1; pn(s,n)
notes(s,"3-4 分钟:\n• 快速过 9 个 — 每个 ~20 秒\n• 让举手: 「最想当哪个? 为什么?」\n• 不深入 — 只是激发兴趣\n• 提醒: 「这只是 9 个! 长大有几千种工作可以选!」")

# 16.6  🔍 SILHOUETTE GUESS GAME — pick 3 jobs from the Cool Jobs Showcase (slide 31)
# Same format on each slide: silhouette/clues on top, ANSWER strip at bottom
# (teacher covers the answer strip with a sticky-note OR just doesn't scroll until students guess)
def silhouette_guess_slide(emoji,job_cn,job_en,color,clue1_cn,clue1_en,clue2_cn,clue2_en,clue3_cn,clue3_en,fact):
    s=ns(); bg(s,CREAM)
    hb(s,"🔍 看图猜职业!  Guess the Job from the Silhouette!",color)
    # LEFT — silhouette area (dark colored box with shadow emoji)
    img=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.30),Inches(3.10))
    img.fill.solid(); img.fill.fore_color.rgb=DARK; img.line.color.rgb=color; img.line.width=Pt(3)
    # Silhouette emoji in semi-faded color
    tb(s,0.4,1.55,4.30,1.5,emoji,sz=140,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.30,4.30,0.40,"⬛ 谁的剪影？",sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.70,4.30,0.30,"Whose silhouette?",sz=11,c=LGRAY,a=PP_ALIGN.CENTER)
    # RIGHT — 3 clues stacked (gradually easier)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(3.10))
    panel.fill.solid(); panel.fill.fore_color.rgb=WHITE; panel.line.color.rgb=color; panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(0.95),Inches(4.85),Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb=color; head.line.fill.background()
    tb(s,5.0,1.03,4.65,0.4,"🔎 3 个提示  3 Clues",sz=14,b=True,c=WHITE)
    clues=[(clue1_cn,clue1_en),(clue2_cn,clue2_en),(clue3_cn,clue3_en)]
    for i,(cn,en) in enumerate(clues):
        y=1.65+i*0.78
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(5.0),Inches(y),Inches(0.40),Inches(0.40))
        badge.fill.solid(); badge.fill.fore_color.rgb=color; badge.line.fill.background()
        tb(s,5.0,y+0.04,0.40,0.32,str(i+1),sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,5.50,y-0.02,4.10,0.45,cn,sz=14,b=True,c=DARK)
        tb(s,5.50,y+0.40,4.10,0.30,en,sz=9,c=GRAY)
    # BOTTOM — guess sentence frame (NO answer shown — teacher reveals verbally)
    rev=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.20),Inches(9.30),Inches(0.95))
    rev.fill.solid(); rev.fill.fore_color.rgb=WARM; rev.line.color.rgb=color; rev.line.width=Pt(2.5)
    tb(s,0.5,4.30,1.6,0.35,"💬 我来猜!",sz=14,b=True,c=color)
    tb(s,2.0,4.27,7.4,0.45,"我猜是 ___ !",sz=22,b=True,c=DARK)
    tb(s,2.0,4.75,7.4,0.30,"I guess it's a ___!",sz=11,c=GRAY)
    return s

# Round 1: 🚀 Astronaut
s=silhouette_guess_slide("🚀","宇航员","Astronaut",NAVY,
    "❤️ 我的兴趣是: 飞 去 很远 很远 的 地方。","❤️ My interest: flying to faraway places.",
    "💪 我擅长: 不晕, 用 工具 在 没有 重力 工作。","💪 I'm good at: working with zero gravity.",
    "🌍 我的 工作 地方 — 没有 空气, 能 看到 整个 地球!","🌍 My workplace: no air, can see the whole Earth!",
    "上过太空 = 全世界 只有 几百 人!")
n+=1; pn(s,n)
notes(s,"猜谜 (1-2 分钟):\n• 让 大孩子 念中文, 小孩子 猜\n• Clue 1 (抽象 — 飞远) → Clue 2 (工具 + 没重力) → Clue 3 (没空气 + 看地球 = 答案)\n• 答案: 宇航员 — 老师 等学生 猜过后 口头揭晓")

# Round 2: 🦒 Vet
s=silhouette_guess_slide("🦒","兽医","Vet (Animal Doctor)",WILD,
    "❤️ 我的兴趣是: 帮 不会 说话 的 朋友。","❤️ My interest: helping friends who can't speak.",
    "💪 我擅长: 听 心跳, 安静 靠近, 不让 病人 害怕。","💪 I'm good at: hearing heartbeats, getting close gently.",
    "🐶 我 给 小狗、小猫、有时 还有 小老虎 打 针。","🐶 I give shots to dogs, cats — sometimes baby tigers!",
    "你家 的宠物 生病 → 找我!")
n+=1; pn(s,n)
notes(s,"猜谜 (1-2 分钟):\n• 大孩子 念, 小孩子 猜\n• Clue 1 (抽象 — 不会说话的朋友) — 学生 可能 猜 婴儿、外国人, 引导他们 想 动物!\n• Clue 2 (听心跳 + 安静靠近) → 像 医生, 但 病人不一样\n• Clue 3 (具体动物名) — 答案 出来\n• 联想: 「你家 有 宠物 吗? 生病时 谁帮忙?」")

# Round 3: 🎮 Game Designer
s=silhouette_guess_slide("🎮","游戏设计师","Game Designer",CITY,
    "❤️ 我的兴趣是: 想 出 好玩 的 故事 和 关卡。","❤️ My interest: making up fun stories and levels.",
    "💪 我擅长: 画 角色, 让 角色 跳过 100 关。","💪 I'm good at: drawing characters, making them jump 100 levels.",
    "🎮 你玩 的 Minecraft / Roblox / 马里奥 = 我们 做 的!","🎮 Minecraft / Roblox / Mario — all made by us!",
    "30 年前 这个工作 还没有! 现在最火!")
n+=1; pn(s,n)
notes(s,"猜谜 (1-2 分钟):\n• 大孩子 念, 小孩子 猜\n• Clue 1 (抽象 — 故事 + 关卡) — 不直接 说 「游戏」, 让 学生 自己 联想\n• Clue 2 (画角色 + 跳关) → 比较 明显 但还没说 「游戏」\n• Clue 3 (具体游戏名) — 学生 全炸\n• 收尾: 「30 年前 没这工作 — 兴趣 + 时代 = 新职业!」")

# 16.7  SPOTLIGHT — 机器人工程师 (Robotics Engineer)
ROBOT=RGBColor(0x55,0x6B,0x83)  # cool steel
s=guess_clue_slide("🤖","Robotics Engineer",ROBOT,
    [("❤️ 我的兴趣是: 让 没有 生命 的 东西 自己 动 起来。","❤️ My interest: making lifeless things move by themselves."),
     ("💪 我擅长: 写 程序, 画 图, 解决 难题。","💪 I'm good at: coding, drawing plans, solving problems."),
     ("🤖 我做的 「员工」 不睡觉、不吃饭, 去过 工厂、医院、太空!","🤖 My 'workers' don't sleep or eat — been to factories, hospitals, space!")],
    "我猜是 ___ !",
    "I guess it's a ___!")
n+=1; pn(s,n)
notes(s,"猜谜 (2-3 分钟):\n• 慢慢念 3 条线索, 一条比一条具体\n• Clue 1 (最模糊): 「不睡觉不吃饭的员工」 — 学生想 (动物? 机器?)\n• Clue 2 (中等): 「画图、写步骤」 — 引出 design + programming\n• Clue 3 (最具体): 工厂/医院/太空 都有 — 引出 robots\n• 让 2-3 个学生猜, 不直接给答案 — 翻页揭晓!")

# Robotics — answer reveal
s=reveal_job_slide("🤖","机器人工程师","Robotics Engineer",ROBOT,
    "设计、编程、造机器人","design, code, build robots",
    "工厂 / 实验室 / 公司","factory / lab / company",
    video_url="搜: 'robotics engineer for kids' (BrainPOP / Crash Course Kids)")
n+=1; pn(s,n)
notes(s,"揭晓 机器人工程师 (1-2 分钟):\n• 全班一起喊: 「机器人工程师!」\n• 看 1 分钟视频 (老师提前找好链接)\n• 复习: 做什么? 在哪里?\n• 句型: 「我是机器人工程师, 我设计能 ___ 的机器人。」")

# 17. STEP 3 — Career Charades
s=ns(); bg(s,CREAM); hb(s,"🎭 演一演!  Career Charades",CHEF)
tb(s,0.4,0.95,9.2,0.45,"上台演职业 — 不能说话!",sz=22,b=True,c=CHEF,a=PP_ALIGN.CENTER)
tb(s,0.4,1.40,9.2,0.30,"Act out a job — no talking!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
# 3 instructional cards
steps=[("1️⃣","抽一张卡片","Pick a card",CHEF),
       ("2️⃣","只能演, 不能说","Act only — no words!",DOC),
       ("3️⃣","大家一起猜","Class guesses!",HELP)]
for i,(em,cn,en,c) in enumerate(steps):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(3.0),Inches(2.30))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.05,2.9,0.7,em,sz=42,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.85,2.9,0.40,cn,sz=16,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.25,2.9,0.30,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.9,0.40,"⏱️ 30 秒",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.40,
    "你是 ___ 吗？  我是 ___ 。",
    "Are you a ___? I am a ___.")
n+=1; pn(s,n)
notes(s,"8 分钟 — 道具准备 (低投入):\n• 把 5 个职业名写在小纸条上, 折起来。\n• 选 4-5 个志愿者上台演 (30 秒), 不能说话, 不能写。\n• 全班用句型猜: 「你是 医生 吗?」\n• 演完, 演员说: 「我是 ___ 。」\n• K 学生: 简单动作 + 老师帮忙\n• G1-3: 加道具 (假装) 让动作更明显")

# 17.5 — CHARADES CARDS: FRONT (pictures) — print + cut out for students
charades_jobs=[
    ("🩺","医生","Doctor",DOC),
    ("📚","老师","Teacher",TEACH),
    ("👮","警察","Police",POLICE),
    ("👨‍🍳","厨师","Chef",CHEF),
    ("👷","工程师","Engineer",ENG),
    ("🦒","兽医","Vet",WILD),
    ("🚒","消防员","Firefighter",DOC),
    ("✈️","飞行员","Pilot",POLICE),
]
# FRONT: pictures only (4 cols × 2 rows)
s=ns(); bg(s,WHITE); hb(s,"🎴 演一演 卡片 · 正面 (Front) — 老师打印, 剪开",CHEF)
tb(s,0.4,0.85,9.2,0.30,"打印这一页 + 下一页 (背面), 双面打印, 沿虚线剪开 — 学生抽卡演!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,1.15,9.2,0.24,"Print this page + next page (back), double-sided. Cut along dashed lines. Students draw a card to act out!",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
for i,(em,job_cn,job_en,c) in enumerate(charades_jobs):
    col=i%4; row=i//4
    x=0.40+col*2.35; y=1.55+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(1.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.30,2.10,1.30,em,sz=80,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"打印说明:\n• 这一页 = 卡片正面 (图片)\n• 下一页 = 卡片背面 (中文 + English)\n• 用 厚纸 (cardstock) 双面打印\n• 沿卡片边线剪开 — 8 张演一演卡片\n• 给每个上台的学生抽 1 张")

# 17.6 — CHARADES CARDS: BACK (CN + EN text) — mirrored layout for double-sided alignment
# When paper is flipped along long edge, back card N aligns with front card N
s=ns(); bg(s,WHITE); hb(s,"🎴 演一演 卡片 · 背面 (Back) — 中文 + English",CHEF)
tb(s,0.4,0.85,9.2,0.30,"卡片背面 — 双面打印时, 这一页与正面对齐",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,1.15,9.2,0.24,"Card backs — when printed double-sided, these align with the fronts opposite",sz=9,c=LGRAY,a=PP_ALIGN.CENTER)
# Mirrored: front order [1,2,3,4 / 5,6,7,8] → back order [4,3,2,1 / 8,7,6,5]
for i,(em,job_cn,job_en,c) in enumerate(charades_jobs):
    # Position card i at MIRRORED column position
    row=i//4
    front_col=i%4
    back_col=3-front_col   # mirror horizontally
    x=0.40+back_col*2.35; y=1.55+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(1.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.color.rgb=c; sh.line.width=Pt(3)
    tb(s,x+0.05,y+0.50,2.10,0.55,job_cn,sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.05,2.10,0.35,job_en,sz=14,b=True,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"卡片背面打印说明:\n• 背面是镜像布局 — 双面打印翻页后, 每张卡片正面和背面对齐\n• 打印设置: 双面打印, 沿长边翻转\n• 颜色背景 = 易于学生快速识别难度\n• 学生抽卡看背面 (答案), 上台演正面 (图片) 提示给观众")

# 18. STEP 3.5 INTRO — 神秘职业! (improved cards with teaser hints)
s=ns(); bg(s,CREAM); hb(s,"✨ 神秘职业  Mystery Jobs!",ENV)
tb(s,0.4,0.85,9.2,0.50,"我再跟你介绍三个职业。",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.40,9.2,0.30,"3 more jobs — have you heard of them?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
mysteries=[
    ("🌍","Mystery 1","和地球有关的工作","About our Earth",      "💧 水 · 空气 · 树",ENV),
    ("🦁","Mystery 2","和动物有关的工作","About wild animals",   "🦒 救动物 · 看动物",WILD),
    ("🏙️","Mystery 3","和我们的「家」有关","About where we live","🗺️ 街道 · 公园 · 路",CITY),
]
for i,(em,tag,hint_cn,hint_en,domain,c) in enumerate(mysteries):
    x=0.4+i*3.15
    # Card box
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(3.0),Inches(2.95))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(3)
    # Tag badge at top
    badge=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.65),Inches(1.95),Inches(1.70),Inches(0.36))
    badge.fill.solid(); badge.fill.fore_color.rgb=c; badge.line.fill.background()
    tb(s,x+0.65,2.00,1.70,0.30,tag,sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # Big emoji
    tb(s,x+0.05,2.40,2.9,0.85,em,sz=58,a=PP_ALIGN.CENTER)
    # Teaser hint (CN bold, EN italics gray)
    tb(s,x+0.05,3.35,2.9,0.40,hint_cn,sz=15,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.78,2.9,0.30,hint_en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    # Domain keywords (separator + small text)
    sep=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+0.50),Inches(4.18),Inches(2.0),Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb=c; sep.line.fill.background()
    tb(s,x+0.05,4.30,2.9,0.40,domain,sz=11,c=DARK,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.95,
    "我猜是 ___ 的工作?",
    "I guess this is a ___ job?")
n+=1; pn(s,n)
notes(s,"30 秒过渡:\n• 「现在 — 3 个新职业! 看看你能不能猜出来。」\n• 不要直接说出职业名 — 让学生从「领域提示」里想\n• 我们一个一个看图猜!")

# 19-24. Three MYSTERY JOBS — Q + Reveal pairs

# Mystery 1: 环境工程师
s=mystery_job_q_slide("🌍","Environmental Engineer",ENV,
    "💧 检测水 / 清洁",
    "他们在做什么？","What are they doing?",
    "是在玩水, 还是让水变干净？","Playing with water — or cleaning it?",
    video_url="搜: 'environmental engineer for kids'  (YouTube)")
n+=1; pn(s,n)
notes(s,"3-4 分钟:\n• 投影一张「人们检测水/清洁水」的真实照片 (建议老师贴一张)。\n• 老师追问: 「他们手里拿什么? 那个瓶子是干什么的?」\n• 让学生先猜, 不要直接给答案!\n• 让 2-3 个孩子说: 「我觉得他们在 ___ 。」")

s=mystery_job_label_slide("🌍","环境工程师","Environmental Engineer",
    "保护水和空气的人","helps keep water and air clean",ENV,
    "他们让水变干净, 让空气变好!")
n+=1; pn(s,n)
notes(s,"揭晓 (1-2 分钟):\n• 「答案是 — 环境工程师!」\n• 让全班跟读 3 遍: 「环境工程师 — 保护水和空气的人」\n• 简单解释: 「水脏了 → 他们让水变干净。空气脏了 → 他们检查空气。」")

# Mystery 2: 野生动物保护员
s=mystery_job_q_slide("🦁","Wildlife Protector",WILD,
    "🦒 救动物 / 观察",
    "他们在做什么？","What are they doing?",
    "是在抓动物, 还是保护动物？","Catching animals — or protecting them?",
    video_url="搜: 'wildlife ranger video for kids'  (YouTube / Nat Geo Kids)")
n+=1; pn(s,n)
notes(s,"3-4 分钟:\n• 投影一张「人们在帮助 / 观察动物」的真实照片。\n• 老师追问: 「动物开心吗? 这些人对动物好不好?」\n• 让学生猜, 不直接给答案!\n• 让 2-3 个孩子说: 「我觉得他们在 ___ 。」")

s=mystery_job_label_slide("🦁","野生动物保护员","Wildlife Protector",
    "保护动物的人","protects animals",WILD,
    "他们救受伤的动物, 看着不让坏人伤害动物!")
n+=1; pn(s,n)
notes(s,"揭晓:\n• 跟读: 「野生动物保护员 — 保护动物的人」\n• 解释: 「动物受伤 → 他们救; 有人想伤害动物 → 他们保护。」")

# Mystery 3: 城市规划师
s=mystery_job_q_slide("🏙️","City Planner",CITY,
    "🗺️ 城市地图 / 模型",
    "他们在做什么？","What are they doing?",
    "是在画画, 还是设计城市？","Drawing pictures — or designing a city?",
    video_url="搜: 'what is urban planning for kids'  (YouTube)")
n+=1; pn(s,n)
notes(s,"3-4 分钟:\n• 投影一张「城市地图 / 模型」照片。\n• 老师追问: 「这是地图吗? 他们要建什么?」\n• 让学生猜!\n• 让 2-3 个孩子说: 「我觉得他们在 ___ 。」")

s=mystery_job_label_slide("🏙️","城市规划师","City Planner",
    "设计城市的人","designs better cities",CITY,
    "他们决定哪里建公园, 哪里建路, 让城市更好住!")
n+=1; pn(s,n)
notes(s,"揭晓:\n• 跟读: 「城市规划师 — 设计城市的人」\n• 解释: 「公园在哪里? 路怎么走? 学校建哪里? 都是他们设计的!」")

# 25. MYSTERY QUICK CHECK — Match
s=ns(); bg(s,CREAM); hb(s,"🎯 快速检查!  Quick Check!",GOLD)
tb(s,0.4,0.85,9.2,0.40,"你来回答 — 谁是 ___ 的人？",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"You answer — who is the person who ___?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
checks=[("💧","谁让水变干净？","Who keeps water clean?","🌍","环境工程师",ENV),
        ("🦒","谁保护动物？","Who protects animals?","🦁","野生动物保护员",WILD),
        ("🗺️","谁设计城市？","Who designs cities?","🏙️","城市规划师",CITY)]
for i,(qem,q_cn,q_en,a_em,a_cn,c) in enumerate(checks):
    y=1.65+i*1.05
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.95))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(2)
    tb(s,0.55,y+0.20,0.6,0.55,qem,sz=30,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.10,4.5,0.40,q_cn,sz=15,b=True,c=DARK)
    tb(s,1.20,y+0.50,4.5,0.30,q_en,sz=10,c=GRAY)
    # arrow
    tb(s,5.85,y+0.20,0.5,0.5,"→",sz=24,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,6.40,y+0.20,0.6,0.55,a_em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,7.05,y+0.25,2.5,0.4,a_cn,sz=14,b=True,c=c)
sentence_frame_bar(s,4.85,
    "___ 是 ___ 的人。",
    "A ___ is the person who ___.")
n+=1; pn(s,n)
notes(s,"3 分钟 — 老师指着问题, 全班一起喊答案!\n• 也可以让学生举手 / 用手势指着对应的图。\n• 让 1-2 个学生用句型说: 「环境工程师是让水变干净的人。」")

# 26. STEP 4 INTRO — Who Can Help?
s=ns(); bg(s,CREAM); hb(s,"🆘 谁来帮忙？  Who Can Help?",HELP)
tb(s,0.4,0.95,9.2,0.45,"出大事了! 谁能帮忙？",sz=24,b=True,c=ALERT,a=PP_ALIGN.CENTER)
tb(s,0.4,1.45,9.2,0.30,"There's a problem! Who can help?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,2.00,9.2,0.40,"看 3 个问题, 你来选职业 — A、B 还是 C？",sz=16,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,0.4,2.45,9.2,0.30,"3 problems — pick a job: A, B, or C?",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
# 3 helpers as choices
helpers=[("A","🌍","环境工程师",ENV),
         ("B","🦁","野生动物保护员",WILD),
         ("C","🏙️","城市规划师",CITY)]
for i,(letter,em,cn,c) in enumerate(helpers):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.95),Inches(3.0),Inches(1.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=c; sh.line.width=Pt(2.5)
    pill(s,x+0.10,3.05,0.7,0.30,letter,c,sz=14)
    tb(s,x+0.85,3.00,2.05,0.45,em,sz=30,a=PP_ALIGN.LEFT)
    tb(s,x+0.05,3.65,2.9,0.40,cn,sz=16,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,4.05,2.9,0.30,"Helper "+letter,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.65,
    "我选 ___ 来帮忙。",
    "I choose ___ to help.")
n+=1; pn(s,n)
notes(s,"30 秒 — 介绍 3 个 helper:\n• 「这 3 个就是我们的「super helpers」!」\n• 「下面看 3 个问题, 你来选 A、B 还是 C。」")

# 27-32. Three SCENARIOS — Q + A pairs
help_options=[("🌍","环境工程师","Env Engineer",ENV),
              ("🦁","野生动物保护员","Wildlife",WILD),
              ("🏙️","城市规划师","City Planner",CITY)]

# Scenario 1: River dirty
s=scenario_q_slide("🐟","河水变脏了, 鱼生病了。","The river is dirty — fish are sick.",ALERT,help_options)
n+=1; pn(s,n)
notes(s,"问题 1 (2-3 分钟):\n• 老师读问题, 让学生想 30 秒。\n• 「A、B、C — 你选谁?」让学生举 A/B/C 手指。\n• 不直接给答案!\n• 翻页揭晓 ✓")

s=scenario_a_slide("🐟","河水变脏了, 鱼生病了。","The river is dirty — fish are sick.",
    "🌍","环境工程师","Environmental Engineer",ENV,
    "他们让水变干净, 鱼就好了!","they clean the water and fish get better!")
n+=1; pn(s,n)
notes(s,"揭晓: 环境工程师!\n• 让全班说: 「环境工程师帮助鱼!」\n• 句型: 「环境工程师帮助 ___ 。」")

# Scenario 2: Sea turtle stuck
s=scenario_q_slide("🐢","海龟被垃圾卡住了。","A sea turtle is stuck in trash.",ALERT,help_options)
n+=1; pn(s,n)
notes(s,"问题 2:\n• 「海龟受伤了 — 谁来救?」\n• 让学生举 A/B/C。\n• 翻页揭晓!")

s=scenario_a_slide("🐢","海龟被垃圾卡住了。","A sea turtle is stuck in trash.",
    "🦁","野生动物保护员","Wildlife Protector",WILD,
    "他们救受伤的动物!","they rescue hurt animals!")
n+=1; pn(s,n)
notes(s,"揭晓: 野生动物保护员!\n• 句型: 「野生动物保护员帮助 ___ 。」")

# Scenario 3: City messy
s=scenario_q_slide("🏘️","城市没有公园, 很乱。","The city has no park — it's messy.",ALERT,help_options)
n+=1; pn(s,n)
notes(s,"问题 3:\n• 「城市没有公园 — 小朋友没地方玩!」\n• 让学生举 A/B/C。\n• 翻页揭晓!")

s=scenario_a_slide("🏘️","城市没有公园, 很乱。","The city has no park — it's messy.",
    "🏙️","城市规划师","City Planner",CITY,
    "他们设计公园, 让城市更好!","they design parks and make the city better!")
n+=1; pn(s,n)
notes(s,"揭晓: 城市规划师!\n• 句型: 「城市规划师帮助 ___ 。」")

# 33. STEP 5 — What if we didn't have this job? (open-ended, multiple occupations)
s=ns(); bg(s,CREAM); hb(s,"💭 如果没有...?  What If We Didn't Have...?",NAVY)
tb(s,0.4,0.85,9.2,0.40,"想一想 — 如果没有这个职业, 我们的世界会怎样?",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Think — if there were no ___, what would our world be like?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 8 occupations in 4x2 grid — open-ended (no hint pills)
ifs=[("🩺","医生","Doctor",DOC),
     ("📚","老师","Teacher",TEACH),
     ("👮","警察","Police",POLICE),
     ("🚒","消防员","Firefighter",DOC),
     ("👨‍🍳","厨师","Chef",CHEF),
     ("👷","工程师","Engineer",ENG),
     ("🦒","兽医","Vet",WILD),
     ("🌍","环境工程师","Env Engineer",ENV)]
for i,(em,job_cn,job_en,c) in enumerate(ifs):
    col=i%4; row=i//4
    x=0.4+col*2.35; y=1.70+row*1.40
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.20),Inches(1.30))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(2.5)
    tb(s,x+0.05,y+0.12,2.10,0.55,em,sz=36,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.72,2.10,0.30,job_cn,sz=14,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+1.02,2.10,0.22,job_en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.65,
    "如果没有 ___ , 我们的世界会 ___ 。",
    "Without ___, our world would ___.")
n+=1; pn(s,n)
notes(s,"5-8 分钟 — 开放式讨论 (不给标准答案!):\n• 老师指着一个职业, 问: 「如果没有医生, 会怎样?」\n• 不要给提示 — 让学生自由发挥, 收集所有想法\n• 鼓励多种回答: 「人生病了没人帮」「医院空了」「妈妈不能上班照顾我」 都对!\n• 一行一行地问 (1 分钟/职业)\n• 关键 takeaway: 每个职业都让世界更好 — 我们都需要彼此!\n• 高年级 stretch: 让学生想一个老师没列的职业, 「如果没有 ___」")

# 33.5  JOBS KEEP CHANGING — Session 1 closing slide (3 columns: before / now / future)
s=ns(); bg(s,CREAM); hb(s,"⏳ 工作一直在变!  Jobs Keep Changing!",HELP)
tb(s,0.4,0.85,9.2,0.40,"以前 → 现在 → 未来 — 工作一直在变!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Past → Now → Future — jobs change all the time!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 3 columns
cols=[
    ("📜","以前",BROWN,"Before",[("✉️","送信"),("📚","抄书"),("🐎","养马"),("🕯️","点路灯")]),
    ("🌍","现在",NAVY,"Now",[("🎮","游戏设计师"),("🤖","AI 工程师"),("📺","YouTuber"),("🚁","无人机飞手")]),
    ("🚀","未来?",GOLD,"Future?",[("🛸","火星导游"),("🤖","机器人医生"),("🌌","太空建筑师"),("🧬","基因设计师")]),
]
for i,(em,head_cn,cl,head_en,items) in enumerate(cols):
    x=0.30+i*3.20
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.65),Inches(3.0),Inches(2.85))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=cl; sh.line.width=Pt(2.5)
    tb(s,x+0.10,1.75,2.80,0.45,f"{em} {head_cn}",sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.10,2.20,2.80,0.30,head_en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    for j,(je,jcn) in enumerate(items):
        y=2.55+j*0.42
        tb(s,x+0.20,y,0.5,0.35,je,sz=18,a=PP_ALIGN.CENTER)
        tb(s,x+0.75,y+0.03,2.10,0.32,jcn,sz=12,c=DARK)
# Predict-disappearing strip at bottom
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.65),Inches(9.40),Inches(0.95))
sh.fill.solid(); sh.fill.fore_color.rgb=GOLD; sh.line.fill.background()
tb(s,0.45,4.75,9.10,0.35,"🔮 你来预测: 今天的哪些工作, 以后可能没有了？",sz=15,b=True,c=WHITE)
tb(s,0.45,5.10,9.10,0.30,"Predict: which jobs today might disappear in the future?",sz=10,c=WARM)
tb(s,0.45,5.40,9.10,0.20,"💭 提示: 司机 (自动驾驶) · 收银员 (自助结账) · 邮递员 (无人机) · ...",sz=9,c=WHITE)
n+=1; pn(s,n)
notes(s,"3-4 分钟 — 收尾互动:\n• 关键句: 「以前没有这些工作, 现在有了 — 等你长大, 也会有现在没有的工作!」\n• 让 2-3 个孩子说: 「我猜未来会有 ___ 工作」\n• 反过来问: 「今天哪些工作以后可能消失?」 — 司机、收银员、邮递员、电梯操作员、流水线工人...\n• 不评判 — 想象力是这一节的目标\n• 收尾: 「重要的不是猜对 — 是 学会变, 学会喜欢, 学会帮人! 」")

# 34. SESSION 2 DIVIDER
s=div("Session 2  下午 2:00–2:45","📚 复习 + 我会认 + 我会写  Review · Read · Write",GOLD,"📖"); n+=1; pn(s,n)

# 35. REVIEW — what we learned in Session 1
s=ns(); bg(s,CREAM); hb(s,"🎮 复习游戏  Review · Baamboozle",GOLD)
tb(s,0.4,0.95,9.2,0.45,"我们来玩 复习游戏! Let's play a review game!",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.45,9.2,0.30,"Quick game to recap morning vocabulary + concepts",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Big Baamboozle placeholder area
ph=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(2.10),Inches(8.0),Inches(2.55))
ph.fill.solid(); ph.fill.fore_color.rgb=WARM; ph.line.color.rgb=GOLD; ph.line.width=Pt(3)
tb(s,1.0,2.40,8.0,0.85,"🎮",sz=70,a=PP_ALIGN.CENTER)
tb(s,1.0,3.30,8.0,0.45,"老师在这里 放 Baamboozle 链接",sz=18,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,1.0,3.78,8.0,0.30,"Teacher: paste Baamboozle game link / QR here",sz=11,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1.0,4.15,8.0,0.30,"baamboozle.com",sz=12,b=True,c=NAVY,a=PP_ALIGN.CENTER)
# Bottom: instructions
inst=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.85),Inches(9.2),Inches(0.65))
inst.fill.solid(); inst.fill.fore_color.rgb=NAVY; inst.line.fill.background()
tb(s,0.55,4.92,9.0,0.30,"📋 分组抢答 — 答对得分! Teams take turns answering for points!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.22,9.0,0.24,"题目涵盖: 5 常见职业 / 3 神秘职业 / 谁来帮忙 / 「我想当 ___」",sz=10,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5-10 分钟 复习游戏:\n• 老师课前在 baamboozle.com 创建一个游戏 (或用已有的)\n• 题目涵盖早上学的内容: 5 常见职业 + 3 神秘职业 + 谁来帮忙 + 公式 (兴趣+能力+帮助=职业)\n• 把游戏链接贴在这一页 (老师手动添加)\n• 上课时直接打开链接玩\n• 全班分 2-3 组抢答得分\n• 不会的题 — 老师讲解后给小组机会再答")

# 36. 我会认 — 5 jobs, ONE SLIDE PER WORD (modeled on 世界旅行 word card pattern)
read_words=[
    ("🩺","医生","yī shēng","Doctor",DOC,
        "医生帮助生病的人。",
        "📷 医生 / 听诊器 / 白大褂"),
    ("📚","老师","lǎo shī","Teacher",TEACH,
        "老师在学校教我们读书。",
        "📷 老师 / 黑板 / 教室"),
    ("👮","警察","jǐng chá","Police",POLICE,
        "警察让大家很安全。",
        "📷 警察 / 警车 / 制服"),
    ("👨‍🍳","厨师","chú shī","Chef",CHEF,
        "厨师做的菜很好吃。",
        "📷 厨师 / 厨房 / 高白帽"),
    ("👷","工程师","gōng chéng shī","Engineer",ENG,
        "工程师设计了这座大桥。",
        "📷 工程师 / 安全帽 / 图纸"),
]
for em,cn,py,en,c,sent,img_label in read_words:
    s=ns(); bg(s,CREAM); hb(s,f"👀 我会认 · {cn}  I Can Read",c)
    # Left: big character card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.fill.background()
    tb(s,0.5,1.10,4.3,1.4,cn,sz=72,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.40,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=c,a=PP_ALIGN.CENTER)
    # Right: image placeholder
    ib(s,5.3,1.0,4.4,2.5,img_label)
    # Bottom: example sentence
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=c; sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=c)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    n+=1; pn(s,n)
    notes(s,f"2 分钟 — {cn}:\n• 老师指字, 全班齐读 3 遍\n• 看图: 「这是 ___ , 在做什么?」\n• 读例句, 跟读\n• 抽 1-2 个学生用 {cn} 造一个新句子")

# 37. 我会写 — 医生 (stroke order practice)
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 医生  I Can Write · Doctor",DOC)
tb(s,0.4,0.85,9.2,0.40,"练一练 — 写「医生」!",sz=22,b=True,c=DOC,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 医生 (Doctor)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Big character display + stroke info side
char1=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(2.30),Inches(2.85))
char1.fill.solid(); char1.fill.fore_color.rgb=WHITE
char1.line.color.rgb=DOC; char1.line.width=Pt(3)
tb(s,0.4,1.95,2.30,1.95,"医",sz=130,b=True,c=DOC,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,2.30,0.30,"yī (medical)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
char2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.85),Inches(1.65),Inches(2.30),Inches(2.85))
char2.fill.solid(); char2.fill.fore_color.rgb=WHITE
char2.line.color.rgb=DOC; char2.line.width=Pt(3)
tb(s,2.85,1.95,2.30,1.95,"生",sz=130,b=True,c=DOC,a=PP_ALIGN.CENTER)
tb(s,2.85,4.10,2.30,0.30,"shēng (life/born)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
# Right: hints panel
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel.fill.solid(); panel.fill.fore_color.rgb=WHITE
panel.line.color.rgb=DOC; panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head.fill.solid(); head.fill.fore_color.rgb=DOC; head.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写 How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"1️⃣ 「医」 — 7 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  先 ㄣ, 然后里面",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"2️⃣ 「生」 — 5 笔",sz=14,b=True,c=DARK)
tb(s,5.45,3.40,4.10,0.30,"  从上往下, 从左往右",sz=10,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格练 3 遍",sz=12,b=True,c=DOC)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「医生」! 我想当医生。",
    "I can write 医生! I want to be a doctor.")
n+=1; pn(s,n)
notes(s,"7-8 分钟 — 写字练习:\n• 老师在白板演示笔顺。\n• 学生跟着空写 (右手食指在空中)。\n• 然后在田字格练习本上写 3-5 遍。\n• 老师巡视, 帮助纠正笔顺。\n• 提示: 「医」里面像个人, 「生」像树苗发芽。\n• 完成后画 ✓。")

# 38. 我会写 — 老师 (stroke order practice)
s=ns(); bg(s,CREAM); hb(s,"✏️ 我会写 · 老师  I Can Write · Teacher",TEACH)
tb(s,0.4,0.85,9.2,0.40,"练一练 — 写「老师」!",sz=22,b=True,c=TEACH,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"Practice writing 老师 (Teacher)",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
char3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(2.30),Inches(2.85))
char3.fill.solid(); char3.fill.fore_color.rgb=WHITE
char3.line.color.rgb=TEACH; char3.line.width=Pt(3)
tb(s,0.4,1.95,2.30,1.95,"老",sz=130,b=True,c=TEACH,a=PP_ALIGN.CENTER)
tb(s,0.4,4.10,2.30,0.30,"lǎo (old/respected)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
char4=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.85),Inches(1.65),Inches(2.30),Inches(2.85))
char4.fill.solid(); char4.fill.fore_color.rgb=WHITE
char4.line.color.rgb=TEACH; char4.line.width=Pt(3)
tb(s,2.85,1.95,2.30,1.95,"师",sz=130,b=True,c=TEACH,a=PP_ALIGN.CENTER)
tb(s,2.85,4.10,2.30,0.30,"shī (master/teacher)",sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
panel2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(2.85))
panel2.fill.solid(); panel2.fill.fore_color.rgb=WHITE
panel2.line.color.rgb=TEACH; panel2.line.width=Pt(2.5)
head2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.30),Inches(1.65),Inches(4.30),Inches(0.50))
head2.fill.solid(); head2.fill.fore_color.rgb=TEACH; head2.line.fill.background()
tb(s,5.45,1.72,4.10,0.4,"✏️ 怎么写 How to write",sz=13,b=True,c=WHITE)
tb(s,5.45,2.30,4.10,0.40,"1️⃣ 「老」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,2.65,4.10,0.30,"  从上往下, 像「土」+「ㄜ」",sz=10,c=GRAY)
tb(s,5.45,3.05,4.10,0.40,"2️⃣ 「师」 — 6 笔",sz=14,b=True,c=DARK)
tb(s,5.45,3.40,4.10,0.30,"  左边 ㄐ, 右边「巾」",sz=10,c=GRAY)
tb(s,5.45,3.85,4.10,0.40,"📝 在田字格练 3 遍",sz=12,b=True,c=TEACH)
tb(s,5.45,4.20,4.10,0.30,"Practice 3 times in grid paper",sz=9,c=GRAY)
sentence_frame_bar(s,4.65,
    "我会写「老师」! 老师教我们。",
    "I can write 老师! The teacher teaches us.")
n+=1; pn(s,n)
notes(s,"7-8 分钟:\n• 演示笔顺, 学生跟写。\n• 田字格练 3 遍。\n• 提示: 「老」上半像帽子, 「师」右边像旗子。\n• 完成 → 颁发「我会写」贴纸。\n• 时间够: 让学生用「老师」造句, e.g. 「我的老师叫 ___ 」")

# 39. SESSION 3 DIVIDER
s=div("Session 3  下午 3:00–4:30","🎨 Output: 小本子 + 职业帽子  Booklet + Career Hat",CITY,"💼"); n+=1; pn(s,n)

# 39.5 — TEAM PAPER BRIDGE COMPETITION (replaces individual Try It · Engineer)
s=ns(); bg(s,CREAM); hb(s,"🌉 团队比赛 · 纸桥挑战  Team Bridge Challenge!",ENG)
tb(s,0.4,0.85,9.2,0.40,"分组 — 比一比 — 谁的纸桥最结实!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Teams compete — whose paper bridge is strongest?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Materials + rules card (left)
mat_card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.70),Inches(4.50),Inches(2.50))
mat_card.fill.solid(); mat_card.fill.fore_color.rgb=WARM; mat_card.line.color.rgb=ENG; mat_card.line.width=Pt(2.5)
tb(s,0.55,1.78,4.30,0.32,"📋 材料 + 规则  Materials + Rules",sz=14,b=True,c=ENG)
tb(s,0.55,2.18,4.30,0.28,"📄 每队 1 张 A4 纸  ·  1 sheet of A4",sz=11,b=True,c=DARK)
tb(s,0.55,2.48,4.30,0.28,"✂️ 1 把剪刀 + 1 卷胶带  ·  Scissors + tape",sz=11,b=True,c=DARK)
tb(s,0.55,2.78,4.30,0.28,"⏱️ 10 分钟造桥  ·  10 min to build",sz=11,b=True,c=DARK)
tb(s,0.55,3.08,4.30,0.28,"📚 桥架在两本书中间 (10 cm)  ·  Span 10 cm",sz=11,b=True,c=DARK)
tb(s,0.55,3.45,4.30,0.32,"🎯 目标: 撑住最多的硬币!",sz=13,b=True,c=ENG)
tb(s,0.55,3.78,4.30,0.24,"Goal: hold the most coins!",sz=10,c=GRAY)
# Engineer 4-step loop (right)
loop_card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.70),Inches(4.50),Inches(2.50))
loop_card.fill.solid(); loop_card.fill.fore_color.rgb=WHITE; loop_card.line.color.rgb=ENG; loop_card.line.width=Pt(2.5)
tb(s,5.25,1.78,4.30,0.32,"🔁 工程师 4 步  Engineer's 4 Steps",sz=14,b=True,c=ENG)
steps=[("1","看","See","观察 the gap"),
       ("2","想","Think","Design the bridge"),
       ("3","做","Make","Build with paper"),
       ("4","改","Improve","Test, fix, retest!")]
for i,(num,cn,en,detail) in enumerate(steps):
    yy=2.20+i*0.46
    pill(s,5.25,yy,0.45,0.36,num,ENG,sz=14)
    tb(s,5.80,yy+0.02,1.0,0.32,cn,sz=16,b=True,c=ENG)
    tb(s,6.85,yy+0.04,1.05,0.30,en,sz=11,b=True,c=DARK)
    tb(s,7.95,yy+0.07,1.65,0.24,detail,sz=9,c=GRAY)
# Competition / scoring (bottom)
score_card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.35),Inches(9.20),Inches(1.0))
score_card.fill.solid(); score_card.fill.fore_color.rgb=ENG; score_card.line.color.rgb=GOLD; score_card.line.width=Pt(2.5)
tb(s,0.55,4.42,9.0,0.32,"🏆 比赛 + 评奖  Test & Award",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,4.78,9.0,0.30,"每队测试 — 撑硬币最多的赢!  也给「最有创意」「最团结合作」奖!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.10,9.0,0.22,"Each team tests — most coins = winner. Also award 'Most Creative' + 'Best Teamwork'!",sz=9,c=GOLD,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"团队比赛 (15-20 分钟):\n• 老师提前准备: A4 纸 (每队 1 张), 剪刀, 胶带, 一些硬币 (or 小石头), 两本厚度相同的书\n• 分组: 4-5 人一队, 共 4-6 队\n• 1 看: 老师摆好两本书 (10 cm gap) — 让学生看「桥」要架的位置\n• 2 想: 每队 2 分钟头脑风暴怎么折/卷/拼 — 不动手!\n• 3 做: 10 分钟造桥 — 老师巡视, 鼓励, 不直接帮\n• 4 改: 第 1 次测试 → 调整 → 第 2 次测试\n• 比赛: 老师给每队的桥逐个加硬币, 直到塌. 撑最多的赢!\n• 颁奖 (3 个): 「最强桥」「最有创意」「最团结合作」 — 每队都有奖\n• 反思: 「桥塌了没关系 — 工程师就是改了再试!」")

# 40. PROJECT 1 — My Dream Job Poster (4 sections aligned to formula)
s=ns(); bg(s,CREAM); hb(s,"🪧 Project 1: 我的梦想职业海报  My Dream Job Poster",CITY)
tb(s,0.4,0.85,9.2,0.36,"画一张海报 — 4 个部分! Draw a poster — 4 parts!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.22,9.2,0.26,"Use the formula: 兴趣 + 能力 + 帮助 = 职业",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# 4 poster sections in 2x2 grid
sections=[
    ("💼","我的梦想职业","Dream Job","我想当 ___ 。","I want to be ___",NAVY),
    ("❤️","我的兴趣","My Interest","我喜欢 ___ 。","I love ___",CITY),
    ("💪","我的能力","My Skill","我会 / 我在练 ___ 。","I can / I'm working on ___",HELP),
    ("🤝","我帮助谁","Who I Help","这个工作帮助 ___ 。","This job helps ___",DOC),
]
for i,(em,cn,en,task_cn,task_en,c) in enumerate(sections):
    col=i%2; row=i//2
    x=0.4+col*4.65; y=1.55+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.45))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=c; sh.line.width=Pt(2.5)
    tb(s,x+0.10,y+0.10,0.85,0.6,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+1.05,y+0.10,3.45,0.36,cn,sz=15,b=True,c=c)
    tb(s,x+1.05,y+0.45,3.45,0.24,en,sz=10,c=GRAY)
    tb(s,x+1.05,y+0.78,3.45,0.30,task_cn,sz=12,b=True,c=DARK)
    tb(s,x+1.05,y+1.06,3.45,0.24,task_en,sz=9,c=GRAY)
sentence_frame_bar(s,4.75,
    "我想当 ___ , 因为 我喜欢 ___ , 我会 ___ , 我想帮助 ___ 。",
    "I want to be ___ because I love ___, I can ___, and I want to help ___.")
n+=1; pn(s,n)
notes(s,"PROJECT 1 · 20-25 分钟:\n• 老师发 1 张大纸 (A3 / 11×17)\n• 学生 折成 4 格 (或老师提前画好 4 框)\n• 4 部分对应公式:\n  - 第 1 部分: 我的梦想职业 (画 + 写)\n  - 第 2 部分: 我的兴趣 (画喜欢的)\n  - 第 3 部分: 我的能力 (画擅长的, 或正在练习的)\n  - 第 4 部分: 我帮助谁 (画受益的人)\n• 差异化: K 画为主, G1-3 加句子, G4-5 写完整段落\n• 完成后挂在墙上 — 准备 Gallery Walk")

# 41. PROJECT 2 — Dream Job Hat AND Badge
s=ns(); bg(s,CREAM); hb(s,"🎩 Project 2: 梦想职业帽子 + 徽章  Hat + Badge!",GOLD)
tb(s,0.4,0.85,9.2,0.40,"选你的梦想职业 — 做一顶帽子 + 一个徽章!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Pick your dream job — make a hat AND a badge!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Two columns: HAT examples (left), BADGE examples (right)
hat_card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.55),Inches(3.10))
hat_card.fill.solid(); hat_card.fill.fore_color.rgb=WHITE; hat_card.line.color.rgb=GOLD; hat_card.line.width=Pt(2.5)
hat_h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.55),Inches(0.45))
hat_h.fill.solid(); hat_h.fill.fore_color.rgb=GOLD; hat_h.line.fill.background()
tb(s,0.55,1.71,4.30,0.33,"🎩 帽子 Hat — 戴上演一演!",sz=14,b=True,c=WHITE)
hats=[("👨‍⚕️","医生帽","白色 + 红十字"),
      ("👮","警察帽","蓝色 + 警徽"),
      ("👨‍🍳","厨师帽","白色, 高高的"),
      ("👷","工程师帽","黄色 + 安全帽")]
for i,(em,cn,hint) in enumerate(hats):
    col=i%2; row=i//2
    x=0.55+col*2.10; y=2.25+row*1.20
    tb(s,x,y,0.55,0.55,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,x+0.55,y+0.05,1.45,0.30,cn,sz=12,b=True,c=DARK)
    tb(s,x+0.55,y+0.32,1.45,0.20,hint,sz=8,c=GRAY)
# Badge examples (right)
bdg_card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.65),Inches(4.55),Inches(3.10))
bdg_card.fill.solid(); bdg_card.fill.fore_color.rgb=WHITE; bdg_card.line.color.rgb=DOC; bdg_card.line.width=Pt(2.5)
bdg_h=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.65),Inches(4.55),Inches(0.45))
bdg_h.fill.solid(); bdg_h.fill.fore_color.rgb=DOC; bdg_h.line.fill.background()
tb(s,5.20,1.71,4.30,0.33,"🎖️ 徽章 Badge — 别在胸前!",sz=14,b=True,c=WHITE)
badges=[("⚕️","医生徽章","红十字 + 名字"),
        ("🛡️","警察徽章","盾牌 + 编号"),
        ("🍴","厨师徽章","刀叉 + 帽子"),
        ("⚙️","工程师徽章","齿轮 + 工具")]
for i,(em,cn,hint) in enumerate(badges):
    col=i%2; row=i//2
    x=5.20+col*2.10; y=2.25+row*1.20
    tb(s,x,y,0.55,0.55,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,x+0.55,y+0.05,1.45,0.30,cn,sz=12,b=True,c=DARK)
    tb(s,x+0.55,y+0.32,1.45,0.20,hint,sz=8,c=GRAY)
sentence_frame_bar(s,4.85,
    "我是小小 ___ ! 我戴 ___ 帽 + ___ 徽章。",
    "I'm a junior ___! I wear a ___ hat + ___ badge.")
n+=1; pn(s,n)
notes(s,"PROJECT 2 · 25-30 分钟:\n• 学生选一个梦想职业 (可以是上午学过的 8 个, 也可以自己想)\n• 做帽子 (15 分钟): 头围纸条 + 帽子图案/装饰\n• 做徽章 (10 分钟): 圆纸片 + 别针/胶带 + 图案\n• 戴上 + 别上 — 「我是小小 ___」 互相介绍\n• 鼓励创意 — 不一定是上面 4 个示例")

# 42. HAT MATERIALS & HOW-TO
s=ns(); bg(s,CREAM); hb(s,"🛠️ 帽子材料 & 做法  Hat Materials & How-To",GOLD)
tb(s,0.4,0.85,9.2,0.40,"简单 4 步 — 30 分钟做一顶帽子!",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.30,"4 simple steps — make a hat in 30 minutes!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Two columns: materials left, how-to right
mat=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.50),Inches(2.95))
mat.fill.solid(); mat.fill.fore_color.rgb=WHITE
mat.line.color.rgb=CHEF; mat.line.width=Pt(2.5)
matH=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.50),Inches(0.50))
matH.fill.solid(); matH.fill.fore_color.rgb=CHEF; matH.line.fill.background()
tb(s,0.55,1.72,4.30,0.4,"📦 材料 Materials",sz=13,b=True,c=WHITE)
materials=[("📄","彩色纸 / 卡纸","colored paper / cardstock"),
           ("✂️","剪刀, 胶水, 胶带","scissors, glue, tape"),
           ("🖍️","彩笔, 贴纸","markers, stickers"),
           ("🪶","羽毛, 棉花 (装饰)","feathers, cotton (deco)"),
           ("🪡","线 / 松紧带 (头围)","string / elastic band")]
for i,(em,cn,en) in enumerate(materials):
    y=2.30+i*0.42
    tb(s,0.55,y,0.40,0.30,em,sz=14)
    tb(s,0.95,y,3.85,0.30,cn,sz=11,b=True,c=DARK)
    tb(s,0.95,y+0.22,3.85,0.20,en,sz=8,c=GRAY)
# How-to right
how=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.65),Inches(4.50),Inches(2.95))
how.fill.solid(); how.fill.fore_color.rgb=WHITE
how.line.color.rgb=HELP; how.line.width=Pt(2.5)
howH=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.65),Inches(4.50),Inches(0.50))
howH.fill.solid(); howH.fill.fore_color.rgb=HELP; howH.line.fill.background()
tb(s,5.25,1.72,4.30,0.4,"🛠️ 4 步做法 4 Steps",sz=13,b=True,c=WHITE)
howto=[("1️⃣","剪一条头围纸条","Cut a paper headband"),
       ("2️⃣","加图案 (红十字 / 警徽 / 帽尖)","Add design (cross / badge / peak)"),
       ("3️⃣","装饰 (颜色 / 贴纸)","Decorate (colors / stickers)"),
       ("4️⃣","卷起来粘成圈 — 戴上!","Tape into a ring — wear it!")]
for i,(em,cn,en) in enumerate(howto):
    y=2.30+i*0.55
    tb(s,5.25,y,0.45,0.4,em,sz=14,b=True,c=HELP)
    tb(s,5.75,y,3.75,0.35,cn,sz=12,b=True,c=DARK)
    tb(s,5.75,y+0.27,3.75,0.25,en,sz=9,c=GRAY)
sentence_frame_bar(s,4.75,
    "我做 ___ 的帽子。我戴上, 我是小小 ___ !",
    "I made a ___ hat. I wear it — I'm a junior ___!")
n+=1; pn(s,n)
notes(s,"老师准备 (低投入):\n• 提前剪好头围纸条 (省时间)。\n• 准备彩纸 + 装饰品 1 桌一份。\n• 老师演示 1 顶 (e.g. 厨师帽: 白色, 加纸团做高顶)。\n• 学生 25 分钟做完, 5 分钟戴上展示。\n• 鼓励合作 — 同桌互相帮忙。")

# 43. PROJECT 3 — Future Jobs Group Project
s=ns(); bg(s,CREAM); hb(s,"🔮 Project 3: 未来 会消失 的 工作  Future Jobs · Group Project",NAVY)
tb(s,0.4,0.85,9.2,0.40,"分组讨论 — 哪些工作 未来 可能 会消失? 为什么?",sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.28,"Group project — which jobs may disappear in the future? Why?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Left: examples / candidates
ex=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.55),Inches(2.55))
ex.fill.solid(); ex.fill.fore_color.rgb=WHITE; ex.line.color.rgb=NAVY; ex.line.width=Pt(2.5)
exh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.55),Inches(0.45))
exh.fill.solid(); exh.fill.fore_color.rgb=NAVY; exh.line.fill.background()
tb(s,0.55,1.71,4.30,0.33,"💡 例子: 哪些可能消失? Examples",sz=13,b=True,c=WHITE)
candidates=[
    ("📞","电话接线员","Phone operator"),
    ("🚕","出租车司机","Taxi driver (自驾车?)"),
    ("📰","报纸送报员","Newspaper carrier"),
    ("🧾","收银员","Cashier (自助结账?)"),
]
for i,(em,cn,en) in enumerate(candidates):
    y=2.20+i*0.48
    tb(s,0.60,y,0.55,0.40,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,1.20,y+0.04,3.6,0.30,cn,sz=12,b=True,c=DARK)
    tb(s,1.20,y+0.32,3.6,0.20,en,sz=8,c=GRAY)
# Right: task instructions
task=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.65),Inches(4.55),Inches(2.55))
task.fill.solid(); task.fill.fore_color.rgb=WHITE; task.line.color.rgb=HELP; task.line.width=Pt(2.5)
taskh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(1.65),Inches(4.55),Inches(0.45))
taskh.fill.solid(); taskh.fill.fore_color.rgb=HELP; taskh.line.fill.background()
tb(s,5.20,1.71,4.30,0.33,"📋 你的任务  Your Task",sz=13,b=True,c=WHITE)
tasks=[
    ("1️⃣","列 4 个 以上 工作","List at least 4 jobs"),
    ("2️⃣","写 / 说 为什么 会消失","Why might they disappear?"),
    ("3️⃣","画 / 写 你的想法","Draw or write your ideas"),
    ("4️⃣","小组分享","Share with the class"),
]
for i,(em,cn,en) in enumerate(tasks):
    y=2.20+i*0.48
    tb(s,5.25,y,0.45,0.40,em,sz=18,b=True,c=HELP)
    tb(s,5.80,y+0.04,3.7,0.30,cn,sz=12,b=True,c=DARK)
    tb(s,5.80,y+0.32,3.7,0.20,en,sz=8,c=GRAY)
# Differentiation bar at bottom
diff=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.30),Inches(9.2),Inches(1.10))
diff.fill.solid(); diff.fill.fore_color.rgb=GOLD; diff.line.fill.background()
tb(s,0.55,4.36,9.0,0.30,"🎯 分级 Differentiation",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,4.66,9.0,0.30,"📚 高年级 (G3+): 用中文写出 4 个原因 (因为 ___)",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,4.94,9.0,0.30,"🎨 低年级 (K-G2): 口头说原因 + 画图 / 写关键词",sz=11,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,5.20,9.0,0.20,"Higher: write reasons in Chinese · Lower: oral reasons + draw / keywords",sz=8,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"PROJECT 3 · 20-25 分钟 — 分组项目:\n• 4-5 人一组, 4-6 组\n• 老师先讨论 1 个例子: 「电话接线员 — 现在我们用手机, 不需要接线员了」\n• 每组用大纸 / 海报:\n  - 1️⃣ 列 4 个以上可能消失的工作\n  - 2️⃣ 解释 为什么 会消失 (新科技? AI? 自动化?)\n  - 3️⃣ 画 / 写 — 视觉化\n  - 4️⃣ 小组上台分享\n• 差异化:\n  - 高年级: 写完整中文句子 「___ 可能会消失, 因为 ___」\n  - 低年级: 老师帮记录, 学生口头说 + 画图\n• 引导思考: 「不要害怕 — 新工作也会出现! 你长大可能做的工作 现在还没发明!」")

# 44. STEP 3 — Gallery Walk
s=ns(); bg(s,CREAM); hb(s,"🚶 画展!  Gallery Walk!",HELP)
tb(s,0.4,0.85,9.2,0.45,"展示你的工具箱 — 给同学看!",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.30,"Show your toolbox to your friends!",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# K vs G1-3 differentiation
k_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.75),Inches(4.50),Inches(2.85))
k_box.fill.solid(); k_box.fill.fore_color.rgb=WHITE
k_box.line.color.rgb=HELP; k_box.line.width=Pt(2.5)
khead=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.75),Inches(4.50),Inches(0.55))
khead.fill.solid(); khead.fill.fore_color.rgb=HELP; khead.line.fill.background()
tb(s,0.55,1.83,4.30,0.4,"🌱 K — 简单版",sz=14,b=True,c=WHITE)
tb(s,0.55,2.50,4.30,0.55,"我是小小 ___ 。",sz=22,b=True,c=DARK)
tb(s,0.55,3.10,4.30,0.40,"I'm a junior ___.",sz=12,c=GRAY)
tb(s,0.55,3.70,4.30,0.40,"💡 1 句 + 指着工具",sz=12,b=True,c=HELP)
tb(s,0.55,4.10,4.30,0.30,"1 sentence + point at tools",sz=9,c=GRAY)
g_box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.75),Inches(4.50),Inches(2.85))
g_box.fill.solid(); g_box.fill.fore_color.rgb=WHITE
g_box.line.color.rgb=GOLD; g_box.line.width=Pt(2.5)
ghead=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.10),Inches(1.75),Inches(4.50),Inches(0.55))
ghead.fill.solid(); ghead.fill.fore_color.rgb=GOLD; ghead.line.fill.background()
tb(s,5.25,1.83,4.30,0.4,"🌟 G1-3 — 完整版",sz=14,b=True,c=WHITE)
tb(s,5.25,2.45,4.30,0.45,"我是小小 ___ 。",sz=18,b=True,c=DARK)
tb(s,5.25,2.95,4.30,0.45,"我用 ___ 。",sz=18,b=True,c=DARK)
tb(s,5.25,3.45,4.30,0.45,"我帮助 ___ 。",sz=18,b=True,c=DARK)
tb(s,5.25,4.00,4.30,0.40,"💡 3 句 + 演示 1 个工具",sz=12,b=True,c=GOLD)
tb(s,5.25,4.40,4.30,0.30,"3 sentences + demo 1 tool",sz=9,c=GRAY)
n+=1; pn(s,n)
notes(s,"15 分钟 — 玩法:\n• 桌面摆好工具箱, 学生轮流当「展览员」和「客人」。\n• 5 分钟换一组。\n• K: 1 句 + 指。\n• G1-3: 3 句 + 演示 (例如「我用听诊器 — 听一听!」)\n• 老师在旁拍照记录。")

# 44.5 — END-OF-DAY · 听绘本 + Turn & Talk (回到 long-after-day reflection)
s=ns(); bg(s,CREAM); hb(s,"🎬 听绘本 + 想一想  Listen & Reflect",GOLD)
tb(s,0.4,0.85,9.2,0.40,"🐶 再听一次:《长大后你想做什么工作》",sz=18,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.4,1.30,9.2,0.28,"Listen again — 'What Do You Want To Be When You Grow Up?'",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
# Left — video reference card (compact)
vcard=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.65),Inches(4.30),Inches(2.85))
vcard.fill.solid(); vcard.fill.fore_color.rgb=WARM; vcard.line.color.rgb=GOLD; vcard.line.width=Pt(3)
tb(s,0.4,1.80,4.30,0.85,"🐶",sz=80,a=PP_ALIGN.CENTER)
tb(s,0.4,2.75,4.30,0.40,"📺 视频  Video",sz=14,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,0.4,3.15,4.30,0.30,"youtube.com/watch?v=EVFPL_qXChU",sz=10,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,0.4,3.55,4.30,0.30,"⏱️ 3-5 分钟",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,3.95,4.30,0.30,"💭 这次 — 不是 听 小狗 的故事,",sz=11,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.20,4.30,0.30,"是 想 — 我自己 长大 想做 什么!",sz=11,b=True,c=GOLD,a=PP_ALIGN.CENTER)
# Right — Turn & Talk prompt (big, focused)
ttcard=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.65),Inches(4.85),Inches(2.85))
ttcard.fill.solid(); ttcard.fill.fore_color.rgb=NAVY; ttcard.line.color.rgb=GOLD; ttcard.line.width=Pt(3)
tb(s,4.85,1.75,4.85,0.40,"🗣️ Turn & Talk",sz=15,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,4.85,2.20,4.85,0.55,"💼 我 长大 想做 什么 工作?",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,4.85,2.78,4.85,0.30,"What do I want to be when I grow up?",sz=10,c=WARM,a=PP_ALIGN.CENTER)
# Chained sentence frame — pulls together everything from today
chain=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.00),Inches(3.20),Inches(4.55),Inches(1.20))
chain.fill.solid(); chain.fill.fore_color.rgb=WARM; chain.line.color.rgb=GOLD; chain.line.width=Pt(2)
tb(s,5.10,3.25,4.40,0.25,"💬 用 完整句型 (Today's Formula!)",sz=10,b=True,c=NAVY)
tb(s,5.10,3.50,4.40,0.30,"我喜欢 ___ ,",sz=12,b=True,c=DARK)
tb(s,5.10,3.75,4.40,0.30,"我擅长 ___ ,",sz=12,b=True,c=DARK)
tb(s,5.10,4.00,4.40,0.30,"我想帮 ___ ,",sz=12,b=True,c=DARK)
tb(s,5.10,4.25,4.40,0.30,"所以 我想当 ___ !",sz=12,b=True,c=GOLD)
# Bottom — closing line ("灵魂 line")
closing=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.65),Inches(9.2),Inches(0.55))
closing.fill.solid(); closing.fill.fore_color.rgb=GOLD; closing.line.fill.background()
tb(s,0.4,4.72,9.2,0.30,"🌟 每个人 都可以 用 自己的 兴趣 和 能力 帮助 世界!",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.4,5.02,9.2,0.20,"Everyone can use their interests and skills to help the world!",sz=10,c=WARM,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"5-7 分钟 — Day 1 收尾反思:\n\n📺 第一阶段 (3-5 分钟): 再放视频\n• 链接: https://www.youtube.com/watch?v=EVFPL_qXChU\n• 这次 不是 听小狗的故事 — 是 让自己 想 「我」长大 想做 什么\n• 不打断 — 让学生 静静 听 + 想\n\n🗣️ 第二阶段 (Turn & Talk, 2-3 分钟):\n• 同桌轮流 用 chain sentence 说:\n  「我喜欢 ___, 我擅长 ___, 我想帮 ___, 所以我想当 ___ 」\n• 这是 today's formula 的最终应用 — 把 4 轮 Turn & Talk 串起来\n\n🎤 第三阶段 (1-2 分钟): 全班 分享\n• 找 3-4 个学生上台 用 chain sentence 说\n• 用 applause-meter 鼓励\n• 收尾 全班齐喊 闪光金句:\n  「🌟 每个人 都可以 用 自己的兴趣 和 能力 帮助 世界!」\n\n• 这是 Day 1 的 灵魂 — 学生离开教室 时, 应该 带着 这个 mantra")

# 45. MISSION COMPLETE!
s=ns(); bg(s,GOLD)
tb(s,1,0.7,8,0.7,"🏆 任务完成!",sz=46,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.55,8,0.5,"Mission Complete!",sz=22,c=WARM,a=PP_ALIGN.CENTER)
# 5 checkmarks
checks=[("✅","猜猜","Guessed"),
        ("✅","演演","Acted"),
        ("✅","神秘","Mystery"),
        ("✅","帮忙","Helped"),
        ("✅","梦想","Dreamed")]
for i,(em,cn,en) in enumerate(checks):
    x=0.4+i*1.88
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.50),Inches(1.78),Inches(1.50))
    sh.fill.solid(); sh.fill.fore_color.rgb=WHITE
    sh.line.color.rgb=NAVY; sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.60,1.7,0.50,em,sz=30,c=OK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.20,1.7,0.40,cn,sz=14,b=True,c=NAVY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.60,1.7,0.30,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.30,9.2,0.45,"🌟 你是小小职业人! Future Career Hero!",sz=22,b=True,c=NAVY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.80,9.2,0.30,"明天继续 — 探索更多职业!",sz=12,c=DARK,a=PP_ALIGN.CENTER)
n+=1; pn(s,n)
notes(s,"庆祝 (1 分钟):\n• 全班一起喊: 「我是小小职业人!」\n• 颁发徽章 / 贴纸 (在工具箱上贴)\n• 「下次见面 — 介绍更多酷职业!」")

# Save
import os
out=os.path.join(os.path.dirname(__file__),"day1_career.pptx")
prs.save(out)
print(f"Saved {out}  ({n} slides)")
