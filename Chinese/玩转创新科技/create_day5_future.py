#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
玩转 创新 科技 · Day 5: 一 周 大 复习 + 科技 改变 生活 大 卷 轴
3-session classroom deck for K-5 Chinese immersion summer camp.

Structure:
  Session 1 (11:00–11:45) — 复 习 Day 1 + Day 2 + Bamboozle 游 戏
  Session 2 (2:00–2:45)   — 复 习 Day 3 + Day 4 + Bamboozle 游 戏
  Session 3 (3:00–4:30)   — Project: 科技 改变 生活 大 卷 轴 (Timeline Mural)
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = make_presentation()
DAY = FUTURE_PINK
ANCIENT = RGBColor(0xB8, 0x50, 0x42)
n = 0


def arrow(s, x, y, w=0.30, h=0.30, color=DAY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


def recap_placeholder(prs, header_text, header_color, msg_cn, msg_en, hint_cn):
    """Blank placeholder slide for teachers to manually insert previous-day slides."""
    s = ns(prs); bg(s, CREAM); hb(s, header_text, header_color)
    # Big centered message panel
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.00), Inches(1.80), Inches(8.00), Inches(2.50))
    box.fill.solid(); box.fill.fore_color.rgb = WARM
    box.line.color.rgb = header_color; box.line.width = Pt(3)
    tb(s, 1.00, 1.95, 8.00, 0.55, "📌", sz=36, a=PP_ALIGN.CENTER)
    tb(s, 1.00, 2.55, 8.00, 0.55, msg_cn, sz=20, b=True, c=header_color, a=PP_ALIGN.CENTER)
    tb(s, 1.00, 3.20, 8.00, 0.35, msg_en, sz=13, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, 1.00, 3.65, 8.00, 0.55, "(留 白 · 人 工 操 作)", sz=14, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    # Hint at bottom
    hf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.60), Inches(9.20), Inches(0.75))
    hf.fill.solid(); hf.fill.fore_color.rgb = WHITE
    hf.line.color.rgb = header_color; hf.line.width = Pt(2)
    tb(s, 0.55, 4.70, 9.0, 0.32, hint_cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 5.05, 9.0, 0.25, "Teachers: leave blank · manually copy + paste key slides from prior PPT here",
       sz=9, c=GRAY, a=PP_ALIGN.CENTER)
    return s


# ============================================================
# 1 · COVER
# ============================================================
cover(prs, 5, "Final Review + Showcase", "一 周 大 复 习 · 科技 卷 轴",
      "🏆 🔥 📜 🤖 🌍", DAY,
      "这 一 周 我 学 到 了 什么?",
      "What did I learn this week?")
n += 1; pn(prs.slides[-1], n)
notes(prs.slides[-1], "Day 5 · 一 周 大 复 习 + 终 极 项目:\n• Session 1: Day 1 (AI) + Day 2 (印刷 / 3D) 复 习 + Bamboozle\n• Session 2: Day 3 (ML) + Day 4 (科技 改变 生活) 复 习 + Bamboozle\n• Session 3: 科技 改变 生活 大 卷 轴 — 全 班 合 作 项目")


# ============================================================
# 2 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1", "🔁 上午 10:00–10:45 / 11:00–11:45  ·  复 习 Day 1 + Day 2 + Bamboozle",
        DAY, "🤖"); n += 1; pn(s, n)


# ============================================================
# 3 · LEARNING GOALS
# ============================================================
s = learning_goals(prs, DAY, [
    ("1️⃣", "回 顾 Day 1 + 2 — 认识 AI + 古代 印刷 / 3D 打印",
     "Review AI + Printing/3D", AI_PURPLE),
    ("2️⃣", "回 顾 Day 3 + 4 — 机器 学习 + 科技 改变 生活",
     "Review ML + Tech changes life", ML_GREEN),
    ("3️⃣", "玩 Bamboozle 复 习 游 戏",
     "Play Bamboozle review games", ORANGE),
    ("4️⃣", "全 班 合作 — 做 「科技 卷 轴」",
     "Class project: Tech Timeline Mural", FUTURE_PINK),
])
n += 1; pn(s, n)


# ============================================================
# 4 · DAY 1 RECAP — placeholder for manual insert
# ============================================================
s = recap_placeholder(prs, "🤖 Day 1 复 习  · 什么 是 AI?", AI_PURPLE,
    "请 挑 选 Day 1 PPT 中 的 重 点 页 面 — 进 行 复 习",
    "Pick key slides from Day 1 PPT — insert here for review",
    "💡 老师: 打 开 day1_ai.pptx → 复 制 重 点 页 → 粘 贴 到 这 一 页 之 后")
n += 1; pn(s, n)
notes(s, "Day 1 复 习 (3-5 分钟):\n• 老师 手 动 挑 选 Day 1 PPT 中 的 重 点 页 (例: AI 是 什么 / 家 里 有 哪 些 AI / AI 主 人 规 则)\n• 把 那 些 页 复 制 + 粘 贴 到 这 一 页 之 后\n• 这 页 留 白 — 不 用 修 改 文 字, 只 是 「插 入 位 置」 提 醒")


# ============================================================
# 5 · DAY 2 RECAP — placeholder for manual insert
# ============================================================
s = recap_placeholder(prs, "🖨️ Day 2 复 习  · 古 代 印 刷 + 3D 打 印", PRINT_ORANGE,
    "请 挑 选 Day 2 PPT 中 的 重 点 页 面 — 进 行 复 习",
    "Pick key slides from Day 2 PPT — insert here for review",
    "💡 老师: 打 开 day2_constellations / 3D printing PPT → 复 制 重 点 页 → 粘 贴 到 这 一 页 之 后")
n += 1; pn(s, n)
notes(s, "Day 2 复 习 (3-5 分钟):\n• 老师 手 动 挑 选 Day 2 PPT 中 的 重 点 页 (例: 毕 昇 + 活 字 印 刷 / 蔡 伦 + 造 纸 / 3D 打 印 是 什 么 / 牙 医 用 3D 打 印)\n• 复 制 + 粘 贴 到 这 一 页 之 后")


# ============================================================
# 6 · BAMBOOZLE Day 1 + 2
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.5, 0.5), (9.0, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.55, 9.4, 0.45, "🎮 GAME TIME!",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.20), Inches(8.8), Inches(2.10))
tt.fill.solid(); tt.fill.fore_color.rgb = AI_PURPLE
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.40, 8.4, 0.85, "Bamboozle · Day 1 + 2 复 习!",
   sz=38, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.30, 8.4, 0.45, "AI + 印刷 / 3D 打印 — 12 道 题!",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.80, 8.4, 0.40, "AI + Printing / 3D Printing — 12 Questions!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

tb(s, 0.3, 3.65, 9.4, 0.95, "🤖  📜  🖨️  ✨  🏆",
   sz=52, a=PP_ALIGN.CENTER)

hf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.55))
hf.fill.solid(); hf.fill.fore_color.rgb = WHITE
hf.line.color.rgb = STAR; hf.line.width = Pt(2)
tb(s, 0.55, 4.92, 9.0, 0.28, "💻 老师: 用 bamboozle_day1_day2_review.csv 导 入 Bamboozle 游 戏!",
   sz=12, b=True, c=AI_PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22, "Teachers: import the CSV to Bamboozle",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "15-20 分钟 Bamboozle 游 戏:\n• 老师 提前 打 开 bamboozle.com 并 上 传 csv\n• 文 件 路径: Chinese/玩转创新科技/bamboozle_day1_day2_review.csv\n• 12 题 涵 盖: AI 是什么 / Siri / AI 会犯错 / 隐私 / 学习方式 / 蔡伦 / 毕昇 / 活字印刷 / 3D 打印\n• 分 2-3 组 比 赛, 答 对 加 分\n• 错 题 老师 立 即 讲 解")


# ============================================================
# 7 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2", "🔁 下午 2:00–2:45  ·  复 习 Day 3 + Day 4 + Bamboozle 游 戏",
        DAY, "🧠"); n += 1; pn(s, n)


# ============================================================
# 8 · DAY 3 RECAP — placeholder for manual insert
# ============================================================
s = recap_placeholder(prs, "🧠 Day 3 复 习  · 机 器 学 习 (Machine Learning)", ML_GREEN,
    "请 挑 选 Day 3 PPT 中 的 重 点 页 面 — 进 行 复 习",
    "Pick key slides from Day 3 PPT — insert here for review",
    "💡 老师: 打 开 day3_ml.pptx → 复 制 重 点 页 (例: ML 六 步 / 兔 子 北 极 熊) → 粘 贴 到 这 一 页 之 后")
n += 1; pn(s, n)
notes(s, "Day 3 复 习 (3-5 分钟):\n• 老师 手 动 挑 选 Day 3 PPT 中 的 重 点 页 (例: 机 器 学 习 六 步 / 兔 子 北 极 熊 / 大 咪 vs 小 咪)\n• 复 制 + 粘 贴 到 这 一 页 之 后")


# ============================================================
# 9 · DAY 4 RECAP — placeholder for manual insert
# ============================================================
s = recap_placeholder(prs, "🌍 Day 4 复 习  · 科技 改变 生活", LIFE_TEAL,
    "请 挑 选 Day 4 PPT 中 的 重 点 页 面 — 进 行 复 习",
    "Pick key slides from Day 4 PPT — insert here for review",
    "💡 老师: 打 开 day4_life.pptx → 复 制 重 点 页 (例: 9 个 发明 时 间 线 / 蔡 伦) → 粘 贴 到 这 一 页 之 后")
n += 1; pn(s, n)
notes(s, "Day 4 复 习 (3-5 分钟):\n• 老师 手 动 挑 选 Day 4 PPT 中 的 重 点 页 (例: 9 个 发明 时 间 线 / 蔡 伦 造 纸 / 四 大 发 明)\n• 复 制 + 粘 贴 到 这 一 页 之 后\n• 提 醒: Session 3 项 目 会 用 到 这 些 — 学 生 要 「选 出 最 重 要 的」")


# ============================================================
# 10 · BAMBOOZLE Day 3 + 4
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.5, 0.5), (9.0, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.55, 9.4, 0.45, "🎮 GAME TIME!",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.20), Inches(8.8), Inches(2.10))
tt.fill.solid(); tt.fill.fore_color.rgb = ML_GREEN
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.40, 8.4, 0.85, "Bamboozle · Day 3 + 4 复 习!",
   sz=38, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.30, 8.4, 0.45, "机器 学习 + 科技 改变 生活 — 13 道 题!",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.80, 8.4, 0.40, "ML + Tech Changes Life — 13 Questions!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

tb(s, 0.3, 3.65, 9.4, 0.95, "🧠  📊  📜  🌍  🏆",
   sz=52, a=PP_ALIGN.CENTER)

hf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.55))
hf.fill.solid(); hf.fill.fore_color.rgb = WHITE
hf.line.color.rgb = STAR; hf.line.width = Pt(2)
tb(s, 0.55, 4.92, 9.0, 0.28, "💻 老师: 用 bamboozle_day3_day4_review.csv 导 入 Bamboozle 游 戏!",
   sz=12, b=True, c=ML_GREEN, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22, "Teachers: import the CSV to Bamboozle",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "15-20 分钟 Bamboozle 游 戏:\n• 文件 路径: Chinese/玩转创新科技/bamboozle_day3_day4_review.csv\n• 13 题 涵 盖: ML 是 什么 / 数据 / 特征 / 兔子 笑 话 / 蔡 伦 / 古 人 写 字 / 科技 = 解决 问题 / 过去 vs 现在\n• 分 组 比 赛 + 错 题 讲 解")


# ============================================================
# 11 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3", "🎨 下午 3:00–4:30  ·  选 一 个 项 目: 📜 时 间 卷 轴  OR  🏙️ 未 来 城 市",
        DAY, "📜"); n += 1; pn(s, n)


# ============================================================
# 12 · STEP 1 — TECH AUCTION (project intro)
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.4, 0.45), (9.1, 0.5), (0.5, 4.85), (9.0, 4.85), (1.2, 0.55), (8.3, 4.95)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.40, 9.4, 0.40, "🏆 Step 1 · 科技 拍 卖 会!",
   sz=18, b=True, c=NEON, a=PP_ALIGN.CENTER)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.95), Inches(8.8), Inches(2.30))
tt.fill.solid(); tt.fill.fore_color.rgb = DAY
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.15, 8.4, 0.85, "🚀 科技 改变 生活 大 卷 轴",
   sz=40, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.10, 8.4, 0.45, "Tech Time Travel Timeline Mural",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.60, 8.4, 0.40, "全 班 合 作 — 从 人 类 历 史 中 挑 选 最 重要 的 发明!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

msg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.45), Inches(8.8), Inches(1.50))
msg.fill.solid(); msg.fill.fore_color.rgb = WHITE
msg.line.color.rgb = STAR; msg.line.width = Pt(3)
tb(s, 0.75, 3.55, 8.5, 0.30, "📣 老师 说:",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.88, 8.5, 0.40, "「卷 轴 空 间 有 限!」",
   sz=20, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 4.28, 8.5, 0.32, "「每 个 时 代 只 能 选 一 个 代 表 发明 — 一 起 来 选!」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.75, 4.65, 8.5, 0.28, "Each era picks ONE invention. Limited scroll space!",
   sz=10, c=GRAY, a=PP_ALIGN.LEFT)

tb(s, 0.3, 5.10, 9.4, 0.30, "⏱️ 15 分 钟 · 拍 卖 + 分 组",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "15 分钟 介 绍 + 分 组:\n• 戏 剧 化 — 像 真 的 拍 卖 师 一 样\n• 强 调 「空 间 有 限」 — 制 造 紧 张 感\n• 然 后 把 全 班 分 成 5 组 (按 时代)\n• 老师 准 备 5 张 大 海报 纸 + 彩 笔 + 贴 纸")


# ============================================================
# 13-17 · 5 ERAS, ONE SLIDE EACH
# ============================================================
era_data = [
    ("Group 1 · 古代 时 代", "Ancient World",
     [("🔥","火","Fire"), ("🛞","轮 子","Wheel"), ("🏹","弓 箭","Bow"),
      ("⛵","帆 船","Sailboat"), ("🌾","农 业","Agriculture")],
     ANCIENT),
    ("Group 2 · 知识 时 代", "Classical & Medieval",
     [("📜","纸","Paper"), ("🧭","指 南 针","Compass"), ("🧮","算 盘","Abacus"),
      ("🏺","陶 器","Pottery"), ("🏗️","水 渠","Aqueduct")],
     PRINT_ORANGE),
    ("Group 3 · 传 播 时 代", "Early Modern",
     [("🖨️","印 刷 术","Printing"), ("🔭","望 远 镜","Telescope"), ("⚙️","蒸汽 机","Steam engine"),
      ("🧪","显 微 镜","Microscope"), ("🌡️","温度 计","Thermometer")],
     ML_GREEN),
    ("Group 4 · 工 业 时 代", "Industrial Age",
     [("☎️","电 话","Phone"), ("💡","电 灯","Light bulb"), ("🚂","火 车","Train"),
      ("🚗","汽 车","Car"), ("✈️","飞 机","Airplane")],
     CYBER),
    ("Group 5 · 数字 时 代", "Digital Age",
     [("💻","电 脑","Computer"), ("🌐","互 联 网","Internet"), ("📱","智 能 手 机","Smartphone"),
      ("🤖","AI","AI"), ("🛰️","GPS","GPS")],
     AI_PURPLE),
]

for group_cn, group_en, candidates, cl in era_data:
    s = ns(prs); bg(s, CREAM, prs); hb(s, f"📜 Step 2 · {group_cn}  ·  {group_en}", cl)

    tb(s, 0.4, 0.85, 9.2, 0.30, "你 们 组 负 责 这 个 时 代 — 从 候 选 中 选 出 「最 重 要 的 一 个」!",
       sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.4, 1.18, 9.2, 0.26, "Your group represents this era — pick the MOST important invention",
       sz=10, c=GRAY, a=PP_ALIGN.CENTER)

    cw_card = 1.78; cgap = 0.10
    ctotal = 5*cw_card + 4*cgap; cstart = (10 - ctotal)/2
    for i, (em, cn, en) in enumerate(candidates):
        x = cstart + i*(cw_card + cgap)
        panel(s, x, 1.65, cw_card, 2.20, cl, fill=WHITE, lw=2.5)
        tb(s, x, 1.80, cw_card, 0.85, em, sz=44, a=PP_ALIGN.CENTER)
        tb(s, x, 2.70, cw_card, 0.40, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
        tb(s, x+0.05, 3.13, cw_card-0.10, 0.28, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
        cb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+cw_card/2-0.18), Inches(3.45), Inches(0.36), Inches(0.36))
        cb.fill.solid(); cb.fill.fore_color.rgb = WHITE
        cb.line.color.rgb = cl; cb.line.width = Pt(2)
        tb(s, x+cw_card/2-0.18, 3.50, 0.36, 0.30, "☐", sz=14, b=True, c=cl, a=PP_ALIGN.CENTER)

    task = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.00), Inches(9.20), Inches(1.55))
    task.fill.solid(); task.fill.fore_color.rgb = cl
    task.line.color.rgb = STAR; task.line.width = Pt(3)
    tb(s, 0.55, 4.08, 9.0, 0.32, "🎯 任 务  Task:",
       sz=13, b=True, c=STAR, a=PP_ALIGN.LEFT)
    tb(s, 0.55, 4.42, 9.0, 0.32, "1.  讨论 — 哪 一 个 最 改变 了 生活? Discuss which is most important.",
       sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
    tb(s, 0.55, 4.75, 9.0, 0.32, "2.  投 票 — 全 组 一 起 选 一 个 (在 □ 里 打 ✓). Vote together.",
       sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
    tb(s, 0.55, 5.08, 9.0, 0.32, "3.  说 出 「为 什么」 — 准 备 给 全 班 解 释. Be ready to explain.",
       sz=11, b=True, c=STAR, a=PP_ALIGN.LEFT)
    n += 1; pn(s, n)
    notes(s, f"5 分钟 小 组 讨论 ({group_cn}):\n• 每 组 4-5 人\n• 让 学 生 各 自 选 一 个 + 说 理 由\n• 然后 小 组 投 票 选 出 「代 表」\n• 强 调: 没 有 「标 准 答 案」 — 重 点 是 讨论 和 思考\n• 老师 走 动 帮 助")


# ============================================================
# 18 · STEP 3 — DISCUSSION GUIDE (3 questions)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "💬 Step 3 · 小 组 讨论  · Discuss in Groups", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "你 们 组 选 好 之 后 — 一 定 要 回 答 这 3 个 问 题!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

discussion_qs = [
    ("1️⃣", "✅", "我们 选 哪 个?",
     "Which invention did we pick?",
     "例如 / Example:  💡 电灯", DAY),
    ("2️⃣", "🌟", "为 什么 最 重要?",
     "Why is it the most important?",
     "句型: We think ___ is the most important. It changed lives because ___.",
     PRINT_ORANGE),
    ("3️⃣", "🤔", "如果 没有 它 怎么 办?",
     "What if it didn't exist?",
     "例如: 没有 电灯 → 晚 上 只 能 用 蜡 烛, 商 店 很 早 关 门 ...",
     AI_PURPLE),
]

for i, (num, em, cn, en, hint, cl) in enumerate(discussion_qs):
    y = 1.30 + i*1.30
    panel(s, 0.40, y, 9.20, 1.20, cl, fill=WHITE, lw=2.5)
    tb(s, 0.55, y+0.20, 0.55, 0.85, num, sz=28, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, 1.20, y+0.18, 0.65, 0.85, em, sz=32, a=PP_ALIGN.LEFT)
    tb(s, 1.95, y+0.15, 7.50, 0.40, cn, sz=17, b=True, c=cl)
    tb(s, 1.95, y+0.55, 7.50, 0.28, en, sz=10, c=GRAY)
    tb(s, 1.95, y+0.85, 7.50, 0.32, hint, sz=11, b=True, c=DARK)
n += 1; pn(s, n)
notes(s, "5-7 分钟:\n• 每 组 必 须 准 备 这 3 个 答 案\n• 老师 走 动 检 查 — 没 想 出 来 的 给 提 示\n• 鼓 励 用 中 文 — 简 单 句 子 就 好\n• Q3 最 有 趣: 「如果 没 有 ___」 — 让 学 生 想象")


# ============================================================
# 19 · STEP 4 — POSTER TEMPLATE
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎨 Step 4 · 做 海 报  · Make Your Poster", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "每 组 一 张 大 海 报 — 中 文 越 多 越 好! 不 会 写 的 就 画!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: Poster template (narrower)
panel(s, 0.40, 1.25, 5.00, 3.95, DAY, fill=WHITE, lw=3)
hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.25), Inches(5.00), Inches(0.50))
hd.fill.solid(); hd.fill.fore_color.rgb = DAY; hd.line.fill.background()
tb(s, 0.40, 1.34, 5.00, 0.35, "📜 海 报 模 板  Poster Template", sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)

tb(s, 0.55, 1.85, 4.70, 0.28, "1.  发明 名 称  + emoji",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 2.12, 4.50, 0.22, "(例: 🔥 火 / Fire)",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

tb(s, 0.55, 2.38, 4.70, 0.28, "2.  画 一 张 大 图",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 2.65, 4.50, 0.22, "(Big illustration)",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

tb(s, 0.55, 2.92, 4.70, 0.28, "3.  为什么 重要?",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.19, 4.50, 0.22, "(Why is it important?)",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

tb(s, 0.55, 3.46, 4.70, 0.28, "4.  如果 没有 它 ...",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.73, 4.50, 0.22, "(What if it didn't exist?)",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

tb(s, 0.55, 4.00, 4.70, 0.28, "5.  影响 评 分  ⭐⭐⭐⭐⭐",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.75, 4.27, 4.50, 0.22, "(Impact rating · 1-5 stars)",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

# Differentiation note inside template
diff = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.60), Inches(4.70), Inches(0.50))
diff.fill.solid(); diff.fill.fore_color.rgb = WARM
diff.line.color.rgb = DAY; diff.line.width = Pt(1.5)
tb(s, 0.60, 4.65, 4.60, 0.22, "💡 中 文 越 多 越 好!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.60, 4.85, 4.60, 0.22, "More Chinese = better!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# RIGHT: Differentiation + Example panels
# Top panel: Two-tier differentiation
diff_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(1.25), Inches(4.05), Inches(1.95))
diff_box.fill.solid(); diff_box.fill.fore_color.rgb = WHITE
diff_box.line.color.rgb = DAY; diff_box.line.width = Pt(2.5)
tb(s, 5.65, 1.32, 3.85, 0.30, "✨ 怎 么 做?  How to make it:",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)

# Drawers row
tb(s, 5.65, 1.65, 0.45, 0.40, "🎨", sz=20, a=PP_ALIGN.LEFT)
tb(s, 6.10, 1.68, 3.45, 0.30, "不 会 写 字 的 — 画 画!",
   sz=12, b=True, c=PRINT_ORANGE, a=PP_ALIGN.LEFT)
tb(s, 6.10, 1.98, 3.45, 0.22, "Can't write yet? Draw pictures!",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

# Writers row
tb(s, 5.65, 2.30, 0.45, 0.40, "✏️", sz=20, a=PP_ALIGN.LEFT)
tb(s, 6.10, 2.33, 3.45, 0.30, "会 写 字 的 — 用 中 文 句 子!",
   sz=12, b=True, c=ML_GREEN, a=PP_ALIGN.LEFT)
tb(s, 6.10, 2.63, 3.45, 0.22, "Can write? Use Chinese sentences/phrases.",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)
tb(s, 6.10, 2.85, 3.45, 0.25, "🌟 句型: 「它 帮 了 ___」「没 有 它 就 ___」",
   sz=10, b=True, c=DARK, a=PP_ALIGN.LEFT)

# Bottom panel: concrete example
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(3.30), Inches(4.05), Inches(1.90))
ex.fill.solid(); ex.fill.fore_color.rgb = STAR
ex.line.color.rgb = DAY; ex.line.width = Pt(2.5)
tb(s, 5.65, 3.38, 3.85, 0.28, "📝 举 个 例 子 · Example: 🧼 肥 皂",
   sz=12, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 5.65, 3.68, 3.85, 0.25, "3.  它 帮 我 们 把 手 洗 干 净.",
   sz=11, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 5.65, 3.95, 3.85, 0.30, "4.  如 果 没 有 它 ...",
   sz=11, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 5.80, 4.25, 3.70, 0.25, "👉 人 们 会 经 常 生 病!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.80, 4.52, 3.70, 0.25, "👉 或 者 ... 把 它 画 出 来!  🤧🤒",
   sz=11, b=True, c=PRINT_ORANGE, a=PP_ALIGN.LEFT)
tb(s, 5.65, 4.82, 3.85, 0.25, "(People would get sick often / Just draw it!)",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

ml = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.30), Inches(9.20), Inches(0.30))
ml.fill.solid(); ml.fill.fore_color.rgb = STAR; ml.line.fill.background()
tb(s, 0.55, 5.32, 9.0, 0.25, "🎨 材 料: 大 海 报 纸 · 彩 笔 · 贴 纸 · 剪刀  ⏱️ 20-25 分钟",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "20-25 分 钟 做 海 报:\n• 老师 准 备: 5 张 大 海 报 纸 (横 向, 长 一 些 — 之 后 要 拼 卷 轴)\n• 彩 笔 / 蜡 笔 / 贴 纸 / 剪刀\n• 差 异 化 教 学:\n  - K-2 (不 会 写 字): 鼓 励 画 画, 老师 帮 写 标 题\n  - 3-5 (会 写 字): 用 中 文 句 子 + 短 语, 越 多 越 好\n• 强 调: 「中 文 越 多 越 好」 — 不 怕 写 错\n• 老师 走 动 提 醒: 「为 什么 这 个 重要?」「没 有 它 会 怎样?」\n• Q4 「如果 没有 它」 是 最 有 创 意 的 — 例 子: 没 肥 皂 → 生 病 / 没 灯 → 摸 黑 / 没 车 → 走 一 整 天")


# ============================================================
# 20 · STEP 5 — ASSEMBLE THE SCROLL
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📜 Step 5 · 拼 成 长 卷  · Assemble the Scroll", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "5 张 海 报 按 时 代 顺 序 拼 起 来 — 一 张 大 时 间 卷 轴!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

eras = [
    ("古 代", "Ancient", "🔥"),
    ("知 识", "Knowledge", "📜"),
    ("传 播", "Communication", "🖨️"),
    ("工 业", "Industrial", "💡"),
    ("数 字", "Digital", "🌐"),
]
era_w = 1.65; era_gap = 0.15
era_total = 5*era_w + 4*era_gap; era_start = (10 - era_total)/2
for i, (cn, en, em) in enumerate(eras):
    x = era_start + i*(era_w + era_gap)
    panel(s, x, 1.55, era_w, 2.50, DAY, fill=WHITE, lw=2.5)
    tb(s, x, 1.70, era_w, 0.80, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 2.55, era_w, 0.40, cn, sz=15, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.95, era_w-0.10, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.35, era_w-0.20, 0.55, "📜 海 报",
       sz=11, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + era_w + 0.01, 2.65, w=0.13, h=0.22, color=DAY)

mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.20), Inches(9.20), Inches(1.35))
mb.fill.solid(); mb.fill.fore_color.rgb = DAY
mb.line.color.rgb = STAR; mb.line.width = Pt(3)
tb(s, 0.55, 4.32, 9.0, 0.40, "🎤 每 组 派 一 个 代 表 — 给 全 班 介 绍!",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.75, 9.0, 0.30, "One representative per group presents to the class!",
   sz=11, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.32, "句型: 「我们 选 ___, 因为 ___. 如果 没 有 ___, 我 们 就 ___.」",
   sz=10, b=True, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 拼 + 5 分钟 介 绍:\n• 老师 把 5 张 海 报 按 时代 顺 序 贴 到 墙 上 / 长 桌 上\n• 形 成 一 个 长 长 的 「时 间 卷 轴」\n• 每 组 选 1 个 代 表 (1 分 钟)\n• 用 句 型: 「我们 选 ___, 因为 ___」")


# ============================================================
# 21 · STEP 6 — CLASS VOTE
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🗳️ Step 6 · 全 班 投 票! · Class Vote", DAY)

tb(s, 0.4, 0.85, 9.2, 0.32, "最 有 趣 的 一 步! 给 你 最 喜欢 的 3 个 发明 贴 贴 纸!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

vote_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.30), Inches(9.20), Inches(1.50))
vote_box.fill.solid(); vote_box.fill.fore_color.rgb = STAR
vote_box.line.color.rgb = DAY; vote_box.line.width = Pt(3)
tb(s, 0.55, 1.42, 9.0, 0.35, "🏆 投 票 规 则  Voting Rules:",
   sz=14, b=True, c=INK, a=PP_ALIGN.LEFT)
vote_rules = [
    "🎯 每 人 拿 3 个 贴 纸 · Each student gets 3 stickers",
    "✨ 贴 在 「改变 世 界 最 多」 的 发明 上 · Stick on the most impactful inventions",
    "📊 数 一 数 — 哪 个 发明 票 最 多? · Count — which one wins?",
]
for i, r in enumerate(vote_rules):
    tb(s, 0.55, 1.80 + i*0.28, 9.0, 0.28, r, sz=12, b=True, c=INK, a=PP_ALIGN.LEFT)

dq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(2.95), Inches(9.20), Inches(1.50))
dq.fill.solid(); dq.fill.fore_color.rgb = DAY
dq.line.color.rgb = STAR; dq.line.width = Pt(3)
tb(s, 0.55, 3.05, 9.0, 0.35, "🤔 大 讨 论  Big Discussion:",
   sz=14, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.45, 9.0, 0.50, "「如果 只 能 留 下 一 个 发明 — 你 会 选 哪 个? 为 什么?」",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.95, 9.0, 0.35, '"If only ONE invention could remain — which would you keep? Why?"',
   sz=12, b=True, c=WARM, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.25, 9.0, 0.25, "(想 想: 没 火 就 没 后 面 的 ... 没 纸 就 没 知识 ... 没 互联网 就 没 现代 ...)",
   sz=10, c=WARM, a=PP_ALIGN.LEFT)

celeb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.60), Inches(9.20), Inches(0.85))
celeb.fill.solid(); celeb.fill.fore_color.rgb = AI_PURPLE
celeb.line.color.rgb = STAR; celeb.line.width = Pt(2.5)
tb(s, 0.55, 4.72, 9.0, 0.35, "🌟 这 一 周 — 你 不 仅 学 了 科技, 你 还 学 会 了 「思 考 + 选 择」!",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.12, 9.0, 0.28, "You learned tech — and how to THINK about it!",
   sz=11, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10-15 分钟 投 票 + 讨论:\n• 老师 准 备 每 人 3 个 贴 纸 (圆 点 / 星 星)\n• 学 生 走 到 卷 轴 前, 把 贴 纸 贴 在 喜 欢 的 发明 上\n• 数 一 数 — 哪 个 票 最 多?\n• 然 后 全 班 讨 论: 「如果 只 能 留 下 一 个 ...」\n• 关 键 引 导:\n  - 没 火 → 后 面 很 多 发明 不 会 出 现\n  - 没 轮子 → 运输 没 法 发展\n  - 没 纸 → 知识 难 保存\n  - 没 互联网 → 现代 生活 完全 不 同\n• 这 是 K-5 高 阶 思 维 的 黄 金 时 刻")


# ============================================================
# 22 · PROJECT 2 DIVIDER — Future City (Alternative)
# ============================================================
s = div(prs, "Project 2", "🏙️ 另 一 个 选 择  ·  全 班 共 同 项 目: 我 们 的 未 来 城 市!",
        DAY, "🏙️"); n += 1; pn(s, n)
notes(s, "另 一 选 择 — 老师 可 以 选 「时间 卷 轴」 (Project 1) 或 「未来 城市」 (Project 2). 都 是 全 班 合 作 + 5 组 分 工.")


# ============================================================
# 23 · PROJECT 2 INTRO — 未来城市 + Photo
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🏙️ Project 2 · 我们 的 未 来 城 市  ·  Our Future City", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "全 班 合 作 — 把 我 们 学 的 科技 用 起 来, 建 一 座 「未 来 城 市」!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: image
img_path = os.path.join(os.path.dirname(__file__), "pics", "future_city.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.40), Inches(1.25), width=Inches(5.40), height=Inches(3.60))
else:
    panel(s, 0.40, 1.25, 5.40, 3.60, DAY, fill=WARM, lw=2)
    tb(s, 0.40, 2.60, 5.40, 0.40, "[未来城市 成 品 照 片]", sz=14, b=True, c=GRAY, a=PP_ALIGN.CENTER)

tb(s, 0.40, 4.95, 5.40, 0.30, "📸 参 考 成 品: 5 个 区 域 拼 成 一 座 城",
   sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

# RIGHT: Concept overview
panel(s, 5.95, 1.25, 3.65, 3.95, DAY, fill=WHITE, lw=2.5)
tb(s, 6.05, 1.35, 3.45, 0.32, "💡 项目 概 念  Concept:",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)

concept_pts = [
    ("🤝", "全 班 合 作", "Whole class"),
    ("5️⃣", "分 成 5 个 区 域", "5 zones"),
    ("🔨", "每 组 做 立 体 模 型", "3D models"),
    ("🎤", "最 后 一 起 展 示!", "Final showcase"),
]
for i, (em, cn, en) in enumerate(concept_pts):
    y = 1.75 + i*0.78
    tb(s, 6.05, y, 0.50, 0.40, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, 6.55, y+0.02, 3.00, 0.30, cn, sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, 6.55, y+0.32, 3.00, 0.24, en, sz=9, c=GRAY, a=PP_ALIGN.LEFT)

ml = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.30), Inches(9.20), Inches(0.30))
ml.fill.solid(); ml.fill.fore_color.rgb = STAR; ml.line.fill.background()
tb(s, 0.55, 5.32, 9.0, 0.25, "⏱️ 60-75 分 钟 · 建 设 + 展 示  ·  Build + Showcase",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "项目 介 绍:\n• 给 学 生 看 成 品 照 片 — 制 造 兴 奋 感\n• 强 调: 「这 是 我 们 今 天 要 做 的!」\n• 5 个 区 拼 在 一 起 = 一 座 城\n• 用 一 周 学 的 科技 (AI / 机器人 / 3D 打印 / 互联网) 想 象 未 来")


# ============================================================
# 24 · PROJECT 2 — 5 GROUPS (zone assignments)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗺️ 5 个 区 域  ·  5 Zones · 你 们 组 是 ...?", DAY)

tb(s, 0.4, 0.85, 9.2, 0.28, "每 组 一 个 区 域 — 设 计 3-4 个 「未 来 科技 产 品」!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

zones = [
    ("Group 1", "🏥 医 疗 区", "Medical Zone",
     ["未来 医 院", "AI 医 生", "自动 手 术 机器人", "健康 扫 描 站"],
     PRINT_ORANGE),
    ("Group 2", "🏫 学 校 区", "School Zone",
     ["未来 学 校", "AI 老 师 助 手", "自动 翻 译 教 室", "VR 学 习 空 间"],
     AI_PURPLE),
    ("Group 3", "🚗 交 通 区", "Traffic Zone",
     ["未来 交 通", "飞 行 汽 车", "自动 公 交", "智 能 红 绿 灯"],
     CYBER),
    ("Group 4", "🏠 家 庭 区", "Home Zone",
     ["未来 家 庭", "做 饭 机 器 人", "自动 清 洁", "智 能 宠 物 助 手"],
     FUTURE_PINK),
    ("Group 5", "🌳 环 保 区", "Eco Zone",
     ["未来 环 保", "垃 圾 分 类 机 器 人", "空 气 清 洁 塔", "太 阳 能 城 市"],
     ML_GREEN),
]

zw = 1.78; zgap = 0.10
ztotal = 5*zw + 4*zgap; zstart = (10 - ztotal)/2
for i, (g, cn, en, items, cl) in enumerate(zones):
    x = zstart + i*(zw + zgap)
    # Card
    panel(s, x, 1.25, zw, 4.00, cl, fill=WHITE, lw=2.5)
    # Header strip
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.25), Inches(zw), Inches(0.85))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.30, zw, 0.28, g, sz=10, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, x, 1.55, zw, 0.40, cn, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 1.95, zw-0.10, 0.20, en, sz=8, c=WARM, a=PP_ALIGN.CENTER)
    # Items
    for j, item in enumerate(items):
        tb(s, x+0.10, 2.25 + j*0.62, zw-0.20, 0.35, f"✦ {item}",
           sz=10, b=True, c=cl, a=PP_ALIGN.LEFT)

# Bottom: presentation prompt
mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.30), Inches(9.20), Inches(0.30))
mb.fill.solid(); mb.fill.fore_color.rgb = STAR; mb.line.fill.background()
tb(s, 0.55, 5.32, 9.0, 0.25, "💡 一 起 想: 在 未 来, 这 个 区 应 该 长 什 么 样? · What does this zone look like in the future?",
   sz=10, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "分 组 + 介 绍 (5-10 分 钟):\n• 老师 随 机 分 5 组 (或 学 生 选)\n• 每 组 看 自 己 的 卡 片\n• 鼓 励: 不 用 全 部 做 — 选 2-3 个 重 点 做 立 体 模 型\n• 提 醒: 用 一 周 学 的 — AI / 机器人 / 3D 打印 / 自动化")


# ============================================================
# 25 · PROJECT 2 — MATERIALS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎨 材 料 + 工 具  ·  Materials & Tools", DAY)

tb(s, 0.4, 0.85, 9.2, 0.28, "简 单 版 — 这 些 就 够 啦!  Simple version — these are enough!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 categories
cats = [
    ("🧱 结 构", "Structure", ["纸 盒 (cardboard)", "LEGO 积 木", "magnetic tiles", "纸 板 管"],
     ANCIENT),
    ("🤖 科技 感", "Tech feel", ["黏 土 (clay)", "锡 纸 (foil)", "吸 管 (straws)", "pipe cleaners"],
     AI_PURPLE),
    ("♻️ 回 收 物", "Recycled", ["瓶 盖", "酸 奶 杯", "鸡 蛋 盒", "小 纸 箱"],
     ML_GREEN),
    ("🎨 装 饰", "Decoration", ["彩 笔 / 蜡 笔", "贴 纸", "googly eyes", "胶 水 + 剪刀"],
     PRINT_ORANGE),
]
cw = 2.20; cgap = 0.10
ctotal = 4*cw + 3*cgap; cstart = (10 - ctotal)/2
for i, (cn_label, en_label, items, cl) in enumerate(cats):
    x = cstart + i*(cw + cgap)
    panel(s, x, 1.20, cw, 3.10, cl, fill=WHITE, lw=2.5)
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.20), Inches(cw), Inches(0.65))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.28, cw, 0.32, cn_label, sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x, 1.58, cw, 0.22, en_label, sz=9, c=WARM, a=PP_ALIGN.CENTER)
    for j, it in enumerate(items):
        tb(s, x+0.12, 1.98 + j*0.50, cw-0.20, 0.32, f"· {it}",
           sz=10, b=True, c=DARK, a=PP_ALIGN.LEFT)

# Bottom: mystery box tip
mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.45), Inches(9.20), Inches(1.00))
mb.fill.solid(); mb.fill.fore_color.rgb = STAR
mb.line.color.rgb = DAY; mb.line.width = Pt(2.5)
tb(s, 0.55, 4.52, 9.0, 0.32, "🎁 神 秘 创 意 盒  Mystery Invention Box!",
   sz=13, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.86, 9.0, 0.28, "老师 准 备 一 箱 随 机 回 收 物 — 让 学 生 自 由 挑! 越 奇 怪 越 有 创 意!",
   sz=11, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.16, 9.0, 0.22, "Random recycled materials in one box · kids dig & invent!",
   sz=9, c=INK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "材料 (最 简 单 版):\n• 必 备: cardboard + LEGO + clay + foil + straws + markers + glue + 回 收 物\n• 老师 提 前 准 备 「mystery invention box」 — 一 箱 杂 物 (瓶 盖 / 酸 奶 杯 / 鸡 蛋 盒 / 小 盒 ...)\n• 每 组 一 个 大 底 板 (poster board) 作 为 「区 域 地 基」\n• 工 具: 剪刀 + 胶 水 (老师 用 热 熔 胶 枪)\n• 小 细 节 — 提 前 打 印: tiny people / cars / trees / 医 院 十 字 标 志")


# ============================================================
# 26 · PROJECT 2 — PRESENTATION FORMAT
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎤 最 后 展 示  ·  Final Presentation", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "每 组 2-3 分 钟 — 介 绍 你 们 的 设 计! 用 这 三 个 句 型 ✨",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 3 sentence frames as big panels
frames = [
    ("1️⃣", "🛠️", "我 们 设 计 的 是 ______.",
     "We designed ___.",
     "例: 「我 们 设 计 的 是 飞 行 汽 车.」", CYBER),
    ("2️⃣", "💖", "它 可 以 帮 助 ______.",
     "It helps ___.",
     "例: 「它 可 以 帮 助 大 家 上 学 不 堵 车.」", PRINT_ORANGE),
    ("3️⃣", "🔧", "它 解 决 的 问 题 是 ______.",
     "It solves the problem of ___.",
     "例: 「它 解 决 的 问 题 是 — 路 上 太 慢, 太 多 车.」", ML_GREEN),
]

for i, (num, em, cn, en, ex, cl) in enumerate(frames):
    y = 1.25 + i*1.10
    panel(s, 0.40, y, 9.20, 1.00, cl, fill=WHITE, lw=2.5)
    tb(s, 0.55, y+0.18, 0.55, 0.65, num, sz=24, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, 1.20, y+0.15, 0.60, 0.70, em, sz=28, a=PP_ALIGN.LEFT)
    tb(s, 1.85, y+0.10, 7.50, 0.40, cn, sz=17, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, 1.85, y+0.48, 7.50, 0.24, en, sz=10, c=GRAY, a=PP_ALIGN.LEFT)
    tb(s, 1.85, y+0.72, 7.50, 0.25, ex, sz=10, b=True, c=DARK, a=PP_ALIGN.LEFT)

# Bottom: showcase
sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.85))
sb.fill.solid(); sb.fill.fore_color.rgb = DAY
sb.line.color.rgb = STAR; sb.line.width = Pt(3)
tb(s, 0.55, 4.78, 9.0, 0.35, "🌆 最 后 — 全 班 拼 起 来 = 「未 来 科技 城 市 展 示」!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.15, 9.0, 0.28, "Combine all 5 zones into one Future Tech City showcase!",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "展 示 (20-30 分 钟):\n• 每 组 2-3 分 钟 — 派 1 个 代 表 或 全 组 一 起 说\n• 三 个 句 型 是 「写 在 黑 板」 的 — 让 学 生 看 着 说\n• 老师 提 问 拓 展: 「这 个 多 久 能 实 现?」「会 不 会 有 问 题?」\n• 最 后 — 5 个 区 拼 在 一 张 大 桌 上 = 城 市!\n• 全 班 合 影 📸\n• 可 选: 用 手 机 录 一 段 「未 来 城 市 导 览」 视 频")


# ============================================================
# 27 · SHARE + CLOSE
# ============================================================
s = share_close(prs, DAY,
    frames_cn=["「这 一 周 我 最 喜 欢 学 的 是 ______」",
               "「最 改变 生 活 的 发明 是 ______」"],
    frames_en="My favorite was ___ · The most impactful invention is ___",
    next_day_cn="✨ 你 已 经 是 「小 创 客」 了! 继续 探 索!",
    next_day_en="You're a Maker now! Keep exploring!",
    next_emoji="🎓")
n += 1; pn(s, n)
notes(s, "10 分钟 收 尾:\n• 每 人 用 句 型 说 一 句\n• 老师 颁 发 「科技 小 工程师」 徽 章 (可 选)\n• 全 班 大 合 影 — 跟 卷 轴 / 城 市 一 起!\n• 一 周 圆 满 结 束 🎉")


out = os.path.join(os.path.dirname(__file__), "day5_future.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
