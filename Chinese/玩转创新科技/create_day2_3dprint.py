#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
玩转 创新 科技 · Day 2: 从 活字 印刷 到 3D 打印  (v2)
副 标题: 从 古代 发明 到 未来 制造

Theme: 中国 古代 创新 (毕昇 活字 印刷) → 现代 科技 (3D 打印) → 学生 hands-on
Pedagogy: storytelling + history + tech + comparison + 2 hands-on projects
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = make_presentation()

# Theme colors
ANCIENT = RGBColor(0xB8, 0x50, 0x42)   # terracotta — ancient/Chinese culture tone
INK_RED = RGBColor(0xC8, 0x25, 0x3E)
MODERN  = CYBER                          # 2A47E0 modern tech blue
DAY     = PRINT_ORANGE                   # day's primary

n = 0


# ============================================================
# 1 · COVER
# ============================================================
s = ns(prs); bg(s, INK, prs)

tb(s, 0.5, 0.35, 9.0, 0.40, "🚀 玩转创新科技  Playing with Innovative Tech",
   sz=16, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.5, 0.85, 9.0, 0.65, "Day 2 · 从活字印刷到 3D 打印",
   sz=28, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.5, 1.55, 9.0, 0.45, "从古代发明到未来制造",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.5, 2.05, 9.0, 0.35, "From Movable Type to 3D Printing · Ancient Invention to Future Making",
   sz=12, c=LGRAY, a=PP_ALIGN.CENTER)

# Two visual placeholders — ancient + modern
photo_slot(s, 1.10, 2.65, 3.50, 2.30,
           "活字印刷 / 毕昇 雕版 图",
           "Movable type / 毕昇 invention image",
           color=ANCIENT)
photo_slot(s, 5.40, 2.65, 3.50, 2.30,
           "3D 打印机 / 立体 模型 图",
           "3D printer / 3D-printed object image",
           color=MODERN)

# Arrow between
tb(s, 4.55, 3.65, 0.95, 0.50, "→", sz=36, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.5, 5.10, 9.0, 0.30, "✨ 古代 → 现代 → 你 的 创造!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Day 2 节奏:\n• Session 1 (45 min): 活字印刷 故事 + 角色 扮演\n• Session 2 (45 min): 3D 打印 + 词汇 + 现场 Demo\n• Session 3 (90 min): 2 个 项目 — DIY 活字印刷 + 3D 涂色\n\n班级: 20+ K-5 immersion 学生\n材料 准备:\n• 一本 厚 书 (用于 hook)\n• 字卡 (我/爱/中/国/学/校 etc.)\n• Foam sheets / carving stamps / sponges (印刷 项目)\n• 提前 3D 打印 一些 玩具 sample\n• 3D 打印机 + 投影 (现场 demo)")

# ============================================================
# 2 · 学习目标 (single-box layout for Day 2)
# ============================================================
s = ns(prs); bg(s, CREAM, prs)
hb(s, "🎯 今天的学习目标  Today's Learning Goals", DAY)
tb(s, 0.40, 0.85, 9.20, 0.32, "上完这节课, 你会……  By the end, you'll be able to…",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# ONE big unified panel
panel(s, 0.40, 1.30, 9.20, 4.10, DAY, fill=WHITE, lw=3)

goals = [
    ("1️⃣", "了解 中国 古代 「活字 印刷」 — 毕昇 的 发明!",
     "Learn about ancient Chinese movable type — Bi Sheng's invention"),
    ("2️⃣", "理解 「印刷」 和 「制造」 都是 把 设计 变 真实",
     "Understand: printing & making = turning design into reality"),
    ("3️⃣", "初步 认识 3D 打印 + 它 能 做 什么",
     "Get to know 3D printing + what it can make"),
    ("4️⃣", "比较 古代 技术 和 现代 科技 — 一样? 不一样?",
     "Compare ancient vs modern tech"),
    ("5️⃣", "动手 做 — 体验 「制造」 的 过程!",
     "Hands-on — experience 'making' yourself!"),
]
for i, (num, cn, en) in enumerate(goals):
    y = 1.50 + i * 0.75
    tb(s, 0.70, y + 0.05, 0.65, 0.50, num, sz=22, b=True, c=DAY)
    tb(s, 1.40, y + 0.05, 8.10, 0.36, cn, sz=14, b=True, c=DARK)
    tb(s, 1.40, y + 0.42, 8.10, 0.28, en, sz=10, c=GRAY)
n += 1; pn(s, n)

# ============================================================
# 3 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1", "📜 上午 11:00–11:45  ·  古代 超级 发明 · 活字 印刷", ANCIENT, "🏛️")
n += 1; pn(s, n)

# ============================================================
# 4 · Part 1 · HOOK — 没有 打印机, 书 怎么 做?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🤔 想一想 · 没有 打印机, 书 怎么 做?", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "老师 拿一本 厚 书 — 你 觉得 古代 人 怎么 做 出 这本 书?",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.40, 1.15, 9.20, 0.22, "Teacher holds a thick book — how was it made BEFORE printers?",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT — book + question photo slot
photo_slot(s, 0.40, 1.50, 4.20, 3.40,
           "一本 古代 / 现代 厚 书 的 照片",
           "Photo of a thick book (ancient or modern)",
           color=DAY)

# RIGHT — 3 guess options
panel(s, 4.80, 1.50, 4.80, 3.40, ANCIENT)
panel_head(s, 4.80, 1.50, 4.80, ANCIENT, "🤷 学生 猜 — 古代 人 怎么 做?", sz=12)
guesses = [
    ("✍️", "一个字 一个字 抄!", "Copy by hand?", CYBER),
    ("🤖", "用 机器?", "Use a machine?", ORANGE),
    ("🎨", "画 出 来?", "Draw it?", PINK),
]
for i, (em, cn, en, cl) in enumerate(guesses):
    y = 2.10 + i * 0.85
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(5.00), Inches(y), Inches(4.40), Inches(0.70))
    chip.fill.solid(); chip.fill.fore_color.rgb = WHITE
    chip.line.color.rgb = cl; chip.line.width = Pt(2)
    tb(s, 5.10, y + 0.10, 0.55, 0.55, em, sz=22)
    tb(s, 5.70, y + 0.08, 3.55, 0.32, cn, sz=14, b=True, c=cl)
    tb(s, 5.70, y + 0.40, 3.55, 0.26, en, sz=9, c=GRAY)

# Bottom — reveal/transition
activity_box(s, 0.40, 5.00, 9.20, 0.55,
             "✋ 举手 投票! 然后 我们 一起 揭晓 — 答案 是 ___ ?",
             "Vote — then we reveal the answer!", color=DAY)
n += 1; pn(s, n)
notes(s, "5 分钟 hook:\n• 老师 真的 拿 一本 书 走进 教室\n• 问: 「这本 书 怎么 做 出 来?」\n• 让 学生 举 手 猜 (3 个 选项 + 自己 说)\n• 不 公布 答案 — 引出 下 一张 slide: 毕昇 + 活字印刷\n\nK-2: 让 他们 摸 一摸 书 — 感受 厚 度\n3-5: 问 「如果 你 是 古代 人, 你 会 怎么 办?」")

# ============================================================
# 5 · 引出 毕昇 + 活字印刷
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📜 答案 揭晓 · 一个 聪明 的 中国 人 — 毕昇!", ANCIENT)

# Big intro card
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.20))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = ANCIENT; intro.line.width = Pt(2.5)
tb(s, 0.55, 1.05, 9.00, 0.45, "🤓 很 久 很 久 以前 (大约 1000 年 前!)",
   sz=14, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.50, 9.00, 0.45, "没有 电脑 + 没有 打印机 — 怎么 办?",
   sz=15, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.95, 9.00, 0.20, "Long, long ago (~1000 years!) — no computers, no printers. What to do?",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: 毕昇 portrait placeholder
photo_slot(s, 0.40, 2.35, 3.50, 2.45,
           "毕昇 画像 / 雕塑 / 古代 画 像",
           "毕昇 (Bi Sheng) portrait or statue",
           color=ANCIENT)

# RIGHT: big idea cards
panel(s, 4.10, 2.35, 5.50, 2.45, ANCIENT)
panel_head(s, 4.10, 2.35, 5.50, ANCIENT, "💡 一个 聪明 的 想法", sz=13)
tb(s, 4.25, 3.00, 5.20, 0.45, "「能 不 能 找 一个 更 快 的 方法?」",
   sz=16, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
tb(s, 4.25, 3.48, 5.20, 0.30, "'Is there a faster way?'", sz=10, italic=False, c=GRAY, a=PP_ALIGN.CENTER) if False else tb(s, 4.25, 3.48, 5.20, 0.30, "'Is there a faster way?'", sz=10, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 4.25, 3.90, 5.20, 0.40, "于是, 毕昇 发明 了 — 活 字 印 刷!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 4.25, 4.32, 5.20, 0.32, "So Bi Sheng invented MOVABLE TYPE printing!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Big quote
tb(s, 0.40, 4.95, 9.20, 0.40, "🇨🇳 这 是 中国 最 重要 的 4 大 发明 之 一!",
   sz=14, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
tb(s, 0.40, 5.35, 9.20, 0.22, "One of China's 4 Great Inventions!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 storytelling:\n• 老师 用 故事 口吻 慢慢 讲\n• 关键 词: 「毕昇」 (bì shēng), 「活字 印刷」, 「1000 年 前」\n• 强调 「中国 古代 发明」 — 文化 自豪感\n• 跟读 3 次 「毕昇」 + 「活字 印刷」\n\n背景 知识 (老师):\n• 毕昇 (大约 970–1051) 北宋 时期\n• 雕版 印刷 在 唐朝 已有, 但 一 版 一 字\n• 毕昇 创新: 单字 可拆 可重 用\n• 4 大 发明: 造纸 / 火药 / 指南针 / 活字 印刷")

# ============================================================
# 6 · 活字印刷 怎么 做 — 4 步
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔨 活字 印刷 4 步  Movable Type · 4 Steps", ANCIENT)
tb(s, 0.40, 0.85, 9.20, 0.30, "毕昇 把 印书 变 简单 — 看 看 怎么 做!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.40, 1.13, 9.20, 0.22, "Bi Sheng made book-printing simple — here's how!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

steps = [
    ("1️⃣", "🔤", "做 一个 一个 字", "Carve each character",
     "把 字 做 成 「小 砖块」", "Like little bricks", ANCIENT),
    ("2️⃣", "🧱", "像 积木 一样 排好", "Line them up like blocks",
     "组成 一句话 / 一首 诗", "Arrange into a sentence", DAY),
    ("3️⃣", "🖌️", "刷 上 墨", "Brush ink on top",
     "黑墨 涂在 字块 上", "Coat with black ink", PINK),
    ("4️⃣", "📄", "压 上 纸 — 印 出 来!", "Press paper on — print!",
     "「印」 就 出 现 在 纸 上!", "And the print appears!", GREEN),
]
card_w = 2.20; gap = 0.10
total = 4 * card_w + 3 * gap; start = (10 - total) / 2
for i, (num, em, cn_t, en_t, cn_d, en_d, cl) in enumerate(steps):
    x = start + i * (card_w + gap)
    panel(s, x, 1.50, card_w, 3.10, cl, lw=2.5)
    # Number badge
    bg_o = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(x + 0.10), Inches(1.62),
                              Inches(0.55), Inches(0.55))
    bg_o.fill.solid(); bg_o.fill.fore_color.rgb = cl; bg_o.line.fill.background()
    tb(s, x + 0.10, 1.68, 0.55, 0.40, num, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Icon
    tb(s, x, 2.30, card_w, 0.80, em, sz=44, a=PP_ALIGN.CENTER)
    # Titles
    tb(s, x + 0.10, 3.15, card_w - 0.20, 0.40, cn_t, sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.10, 3.55, card_w - 0.20, 0.28, en_t, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Description
    tb(s, x + 0.10, 3.90, card_w - 0.20, 0.40, cn_d, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x + 0.10, 4.30, card_w - 0.20, 0.28, en_d, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Arrow between
    if i < 3:
        tb(s, x + card_w, 2.90, gap, 0.40, "→", sz=18, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)

# Bottom photo slot
photo_slot(s, 0.40, 4.80, 9.20, 0.75,
           "活字印刷 完整 过程 流程 图 / 实际 操作 照片",
           "Photo/diagram of the full movable-type process",
           color=ANCIENT)
n += 1; pn(s, n)
notes(s, "8 分钟:\n• 老师 念 4 步, 全班 跟读\n• 用 手势 比划 每 步:\n  1. 假装 刻字 (手 握 笔 / 刀)\n  2. 摆 积木 动作\n  3. 假装 刷 墨\n  4. 拍 一下 桌子 (压 纸 印!)\n\n搞笑 比喻:\n• 字块 = 「小砖块」 / 「小积木」\n• 像 玩 乐高 一样 排 字!\n• 错 了? 换 一个 字 就行!")

# ============================================================
# 7 · 互动问题 — 为什么 活字印刷 聪明?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🧠 互动 · 为什么 活字印刷 这么 聪明?", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "毕昇 的 想法 为什么 厉害? 想 一 想!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 3 question cards with answers
qs = [
    ("❓", "为什么 以前 抄书 很 慢?",
     "Why was copying so slow?",
     "一个字 一个字 写 — 一本书 要 写 好 几 个 月!",
     "Letter by letter — took MONTHS!", ANCIENT),
    ("💡", "为什么 活字印刷 聪明?",
     "Why is movable type smart?",
     "字 可以 「拆 下来 重新 用」 — 不 用 重 刻!",
     "Letters can be REUSED — no need to carve again!", DAY),
    ("🤔", "如果 排 错 字 怎么 办?",
     "What if a letter is wrong?",
     "拿 走 错 字, 换 一个 对 的 — 超 方便!",
     "Just swap it — super easy!", GREEN),
]
card_w = 2.95; gap = 0.12
total = 3 * card_w + 2 * gap; start = (10 - total) / 2
for i, (em, q_cn, q_en, a_cn, a_en, cl) in enumerate(qs):
    x = start + i * (card_w + gap)
    panel(s, x, 1.30, card_w, 3.60, cl, lw=2.5)
    tb(s, x, 1.42, card_w, 0.60, em, sz=34, a=PP_ALIGN.CENTER)
    # Question header
    tb(s, x + 0.10, 2.05, card_w - 0.20, 0.45, q_cn, sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.10, 2.50, card_w - 0.20, 0.30, q_en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Divider
    tb(s, x, 2.85, card_w, 0.18, "─ ─ ─", sz=10, c=cl, a=PP_ALIGN.CENTER)
    # Answer
    tb(s, x + 0.10, 3.10, card_w - 0.20, 0.85, a_cn, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x + 0.10, 3.95, card_w - 0.20, 0.40, a_en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

activity_box(s, 0.40, 5.00, 9.20, 0.55,
             "👯 同桌讨论 1 分钟 — 你 还 想 到 什么 优点?",
             "Talk with partner — what other advantages?", color=DAY)
n += 1; pn(s, n)

# ============================================================
# 8 · 体验 · 我 是 印刷 工人! (Role play)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎭 体验 · 我 是 印刷 工人!  Be a Print Worker!", PINK)

intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.70))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = PINK; intro.line.width = Pt(2)
tb(s, 0.55, 1.02, 9.00, 0.30, "🎯 全班 一起 来 体验 — 用 字卡 排 句子!",
   sz=13, b=True, c=PINK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.32, 9.00, 0.24, "Class activity — arrange word cards into sentences!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: word cards visual
panel(s, 0.40, 1.85, 4.55, 3.20, PINK)
panel_head(s, 0.40, 1.85, 4.55, PINK, "🔤 字 卡  Word Cards", sz=12)
# Display 6 character cards in 2x3 grid
chars = ["我", "爱", "中", "国", "学", "校"]
for i, ch in enumerate(chars):
    col = i % 3; row = i // 3
    x = 0.65 + col * 1.30
    y = 2.50 + row * 1.05
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(1.10), Inches(0.90))
    card.fill.solid(); card.fill.fore_color.rgb = WARM
    card.line.color.rgb = PINK; card.line.width = Pt(2.5)
    tb(s, x, y + 0.15, 1.10, 0.65, ch, sz=32, b=True, c=PINK, a=PP_ALIGN.CENTER)

# RIGHT: example sentences students can make
panel(s, 5.05, 1.85, 4.55, 3.20, GREEN)
panel_head(s, 5.05, 1.85, 4.55, GREEN, "✨ 可以 排 出 什么 句子?", sz=12)
sentences = [
    "「我 爱 中 国」",
    "「我 爱 学 校」",
    "「我 爱 中 国 学 校」",
    "「我 爱 我 国 校 ?」 (排错了!)",
]
for i, sent in enumerate(sentences):
    y = 2.40 + i * 0.55
    cl = GREEN if i < 3 else PINK
    tb(s, 5.20, y, 4.30, 0.45, sent, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)

# Activity callout — reflection (slim compact bar)
race_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.45))
race_bar.fill.solid(); race_bar.fill.fore_color.rgb = PINK; race_bar.line.fill.background()
tb(s, 0.55, 5.16, 9.00, 0.32, "🎯 排错 字? — 拿走 错 字, 换 一个 对 的! 像 毕昇 一样!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 角色 扮演:\n\n准备:\n• 老师 提前 准备 字卡 (大字, 卡纸 or printed) — 每组 一套\n  - 必备字: 我 / 爱 / 中 / 国 / 学 / 校\n  - 加 字: 妈/爸/家/小/朋/友 / 数字 / AI\n• 学生 分 小 组 (3-4 人)\n\n步骤:\n• 老师: 「现在 你们 是 印刷 工人!」\n• 1. 排 「我爱中国」 — 全班 一起\n• 2. 排 「我爱学校」\n• 3. 自由 创作 句子\n• 4. 故意 排错 一个 — 让 学生 「修正」 (体验 「换字」 的 优势)\n\n讨论:\n• 「如果 没有 活字, 错 一个 字 怎么 办?」\n• 答: 「整 块 木板 都 要 重 刻!」 (痛苦!)\n\n分层:\n• K-2: 老师 给 句型, 学生 找字 排\n• 3-5: 自己 想 句子, 找字 排 + 教 同伴")

# ============================================================
# 9 · Part 4 · 连接 — 现在 还能 印 什么?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🌉 古代 印 字 — 现在 印 什么?", PURPLE)

# Top intro card
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.95))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = PURPLE; intro.line.width = Pt(2.5)
tb(s, 0.55, 1.05, 9.00, 0.40, "🤔 古代 人 把 「字」 印 在 纸 上 — 现在 的 机器 还 能 做 什么?",
   sz=14, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.50, 9.00, 0.30, "Ancient people printed words on paper — what can today's machines do?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Big reveal
reveal = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.40), Inches(2.15), Inches(9.20), Inches(2.45))
reveal.fill.solid(); reveal.fill.fore_color.rgb = INK
reveal.line.color.rgb = STAR; reveal.line.width = Pt(3)
tb(s, 0.55, 2.30, 9.00, 0.45, "✨ 现在 不 只 印 字 —",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.85, 9.00, 0.70, "还 能 印 「东西」!",
   sz=36, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 3.65, 9.00, 0.40, "玩具 · 鞋子 · 房子 · 假牙 · 火箭 零件 …",
   sz=16, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.10, 9.00, 0.30, "Today's machines can print OBJECTS — toys, shoes, houses, teeth, rocket parts!",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.40, 9.00, 0.18, "Brace yourself…",
   sz=8, italic=False, c=LGRAY, a=PP_ALIGN.CENTER) if False else tb(s, 0.55, 4.40, 9.00, 0.18, "Brace yourself…", sz=8, c=LGRAY, a=PP_ALIGN.CENTER)

# Bottom teaser
teaser = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.65))
teaser.fill.solid(); teaser.fill.fore_color.rgb = MODERN
teaser.line.color.rgb = STAR; teaser.line.width = Pt(2)
tb(s, 0.55, 4.95, 9.00, 0.30, "🖨️ 下午 见 — 3D 打印!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.25, 9.00, 0.20, "See you this afternoon for 3D printing!",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 衔接:\n• 老师 用 dramatic tone\n• 「你 觉得 ___ 也 能 印 出 来 吗?」\n• 学生 猜: 鞋子? 房子? 玩具?\n• 老师: 「真的! 下午 我们 看 真 的 3D 打印机!」\n\n激发 兴趣 — 让 学生 期待 下午")

# ============================================================
# 10 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2", "🖨️ 下午 2:00–2:45  ·  未来 制造 · 3D 打印", MODERN, "🚀")
n += 1; pn(s, n)

# ============================================================
# 11 · 早上 回顾
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔁 早上 学了 什么?  Morning Recap", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "回想 一下 — 早上 我们 认识 了 什么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

recap = [
    ("📜", "毕昇 是 谁?", "毕昇 (Bi Sheng) — 中国 古代 聪明 人, 发明 了 活字 印刷", ANCIENT),
    ("🔤", "活字印刷 = ?", "把 字 做 成 小 块, 排好 + 刷墨 + 印!", DAY),
    ("💡", "为什么 聪明?", "字 可以 拆下 重用 — 比 抄书 快 100 倍!", GREEN),
    ("🌉", "现在 还能 印 什么?", "不只字 — 还能 印 「东西」! (今天 下午 学!)", PURPLE),
]
for i, (em, q, a, cl) in enumerate(recap):
    y = 1.30 + i * 0.95
    panel(s, 0.50, y, 9.00, 0.80, cl, fill=WHITE, lw=2)
    tb(s, 0.65, y + 0.18, 0.55, 0.50, em, sz=24)
    tb(s, 1.30, y + 0.10, 8.00, 0.32, q, sz=13, b=True, c=cl)
    tb(s, 1.30, y + 0.42, 8.00, 0.32, a, sz=10, b=True, c=DARK)
n += 1; pn(s, n)
notes(s, "3-5 分钟 快速 回顾:\n• 老师 念 问题, 学生 答\n• 全班 喊 「毕昇!」 + 「活字 印刷!」\n• 准备 进入 3D 打印 学习")

# ============================================================
# 12 · 什么 是 3D 打印 — 对比 普通 打印 vs 3D 打印
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "💡 什么 是 3D 打印?  What is 3D Printing?", MODERN)
tb(s, 0.40, 0.85, 9.20, 0.30, "看 对比 — 普通 打印 vs 3D 打印, 有 什么 不一样?",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: 普通 打印 (2D)
panel(s, 0.40, 1.30, 4.55, 3.55, ORANGE)
panel_head(s, 0.40, 1.30, 4.55, ORANGE, "📄 普通 打印  Regular Printing", sz=13)
photo_slot(s, 0.55, 1.85, 4.25, 1.85,
           "普通 打印机 + 平面 纸 张",
           "Regular printer + flat paper",
           color=ORANGE)
tb(s, 0.55, 3.80, 4.25, 0.30, "= 平面 的!  Flat (2D)", sz=14, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.15, 4.25, 0.45, "只 印 字 和 图 (在 纸 上)",
   sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.55, 4.25, 0.28, "Only words and pictures on paper",
   sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# RIGHT: 3D 打印
panel(s, 5.05, 1.30, 4.55, 3.55, MODERN)
panel_head(s, 5.05, 1.30, 4.55, MODERN, "🖨️ 3D 打印  3D Printing", sz=13)
photo_slot(s, 5.20, 1.85, 4.25, 1.85,
           "3D 打印机 + 立体 物 品",
           "3D printer + a 3D object",
           color=MODERN)
tb(s, 5.20, 3.80, 4.25, 0.30, "= 立体 的!  3D (Three-dimensional)", sz=14, b=True, c=MODERN, a=PP_ALIGN.CENTER)
tb(s, 5.20, 4.15, 4.25, 0.45, "能 印 真实 的 「东西」!",
   sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 5.20, 4.55, 4.25, 0.28, "Can print REAL OBJECTS you can hold!",
   sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom analogy
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.60))
band.fill.solid(); band.fill.fore_color.rgb = INK
band.line.color.rgb = STAR; band.line.width = Pt(2)
tb(s, 0.55, 5.02, 9.00, 0.28, "💡 比喻: 像 「挤 奶油」 一样 — 一层 一层 堆 起来!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.30, 9.00, 0.20, "Like squeezing icing — layer by layer until 3D!",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "8 分钟:\n• 老师 用 「挤奶油」 比喻 — 手势 演示\n• 「像 挤 牙膏 一样, 一层 一层 加 上去!」\n• 跟读 关键 词: 「平面」 vs 「立体」\n\n关键 概念:\n• 2D = 长 + 宽 (纸 上)\n• 3D = 长 + 宽 + 高 (可以 拿 起来)")

# ============================================================
# 13-17 · 我会认 (5 vocab cards)
# ============================================================
recognize_words = [
    ("📜", "活字 印刷", "huó zì yìn shuā", "Movable Type Printing",
     "毕昇 发明 了 活字 印刷。", "Bi Sheng invented movable type printing.",
     "活字 印刷 字块 / 排版 照片", ANCIENT),
    ("🖨️", "打印", "dǎ yìn", "Print",
     "我 用 打印机 打印 作 业。", "I use a printer to print homework.",
     "打印机 / 打印 中 的 纸", DAY),
    ("✏️", "设计", "shè jì", "Design",
     "我 设计 了 一 个 机器人。", "I designed a robot.",
     "草图 / 学生 在 画 设计", CYBER),
    ("⚙️", "机器", "jī qì", "Machine",
     "工厂 里 有 很 多 机器。", "Factories have many machines.",
     "工厂 机器 / 齿轮 / 机械", GREEN),
    ("🧊", "模型", "mó xíng", "Model",
     "这 是 太空 站 的 模型。", "This is a space station model.",
     "3D 模型 / 玩具 模型", PURPLE),
]
for em, cn, py, en, ex_cn, ex_en, hint, cl in recognize_words:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)

# ============================================================
# 18-19 · 我会写 (2 writing slides)
# ============================================================
s = vocab_write(prs, DAY, "打印", "Print",
                [("打", "dǎ", "5 笔", "「扌」 旁 + 「丁」"),
                 ("印", "yìn", "5 笔", "左 「卯」 + 右 「卩」 — 像 印章")])
n += 1; pn(s, n)

s = vocab_write(prs, GREEN, "机器", "Machine",
                [("机", "jī", "6 笔", "「木」 旁 + 「几」"),
                 ("器", "qì", "16 笔", "上 「口」 + 下 「口」+「犬」")])
n += 1; pn(s, n)

# ============================================================
# 20 · 3D 打印 怎么 工作 (3 steps)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "⚙️ 3D 打印 怎么 工作?  How Does 3D Printing Work?", MODERN)

steps = [
    ("1️⃣", "✏️", "画 设计 图", "Design",
     "在 电脑 里 画 你 要 的 形 状", CYBER),
    ("2️⃣", "🍦", "加 热 「材料」", "Heat material",
     "塑料 + 高温 → 像 软软 的 奶油", ORANGE),
    ("3️⃣", "🍰", "一层 一层 堆", "Layer by layer",
     "像 在 蛋糕 上 挤 奶油!", PINK),
]
card_w = 2.95; gap = 0.12
total = 3 * card_w + 2 * gap; start = (10 - total) / 2
for i, (num, em, cn_t, en_t, desc, cl) in enumerate(steps):
    x = start + i * (card_w + gap)
    panel(s, x, 1.30, card_w, 3.40, cl, lw=2.5)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.30), Inches(card_w), Inches(0.45))
    head.fill.solid(); head.fill.fore_color.rgb = cl; head.line.fill.background()
    tb(s, x, 1.36, card_w, 0.35, num, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Big icon
    tb(s, x, 1.85, card_w, 0.85, em, sz=54, a=PP_ALIGN.CENTER)
    # Title
    tb(s, x + 0.10, 2.75, card_w - 0.20, 0.40, cn_t, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.10, 3.15, card_w - 0.20, 0.30, en_t, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Description
    tb(s, x + 0.15, 3.55, card_w - 0.30, 1.05, desc, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    if i < 2:
        tb(s, x + card_w, 2.65, gap, 0.50, "→", sz=18, b=True, c=MODERN, a=PP_ALIGN.CENTER)

activity_box(s, 0.40, 4.85, 9.20, 0.70,
             "💡 想一想: 这个 「一层一层 堆」 和 早上 的 「活字 排版」 像 不 像?",
             "Think: layer-by-layer printing vs movable-type stacking — similar?",
             gesture_hint="🤔 都 是 「把 设计 变成 真实」!", color=MODERN)
n += 1; pn(s, n)
notes(s, "5-8 分钟:\n• 用 「挤奶油」 「挤牙膏」 手势 演示\n• 强调 「一层一层」 — 不是 一下 就 出 来\n• 跟读 3 步\n\n衔接 古代:\n• 活字 = 把字 拼 起来 → 印\n• 3D = 把 材料 堆 起来 → 印\n• 都 是 「制造」!")

# ============================================================
# 21 · 3D 打印 能 做 什么 (applications)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🌍 3D 打印 能 做 什么?  What Can 3D Printing Make?", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "你 知道 吗? 这些 东西 都 可以 3D 打印!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

apps = [
    ("🎮", "玩具", "Toys", "小机器人 / 恐龙 / 手办", CYBER),
    ("🦷", "假牙", "Dental", "牙医 用 3D 打印 假牙", PINK),
    ("👟", "鞋子", "Shoes", "Nike / Adidas 都用!", ORANGE),
    ("🏠", "房子", "Houses", "真的 能 打印 房子!", GREEN),
    ("🚀", "火箭 零件", "Rocket parts", "NASA 用 3D 打印!", PURPLE),
    ("🦾", "假肢", "Prosthetics", "帮助 残疾 朋友", MODERN),
]
card_w = 2.95; gap = 0.12; row_gap = 0.15
total = 3 * card_w + 2 * gap; start = (10 - total) / 2
for i, (em, cn_t, en_t, cn_d, cl) in enumerate(apps):
    row = i // 3; col = i % 3
    x = start + col * (card_w + gap)
    y = 1.30 + row * (1.85 + row_gap)
    panel(s, x, y, card_w, 1.85, cl, lw=2)
    tb(s, x + 0.10, y + 0.10, 0.65, 0.65, em, sz=30)
    tb(s, x + 0.75, y + 0.10, card_w - 0.85, 0.40, cn_t, sz=14, b=True, c=cl)
    tb(s, x + 0.75, y + 0.50, card_w - 0.85, 0.28, en_t, sz=9, c=GRAY)
    tb(s, x + 0.15, y + 0.95, card_w - 0.30, 0.80, cn_d, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

activity_box(s, 0.40, 5.20, 9.20, 0.40,
             "🤔 你 想 3D 打印 什么? — 转身 告诉 同 桌!",
             "What would YOU print? Turn and tell your partner!", color=DAY)
n += 1; pn(s, n)
notes(s, "8 分钟 互动:\n• 老师 念 每 个 例子, 学生 看 placeholder 想象 真实 照片\n• 强调 「假牙」 + 「假肢」 — 3D 打印 真的 帮 助 人\n• 让 学生 自由 说 「我 想 打印 ___」\n\n要 准备:\n• 老师 提前 真的 3D 打印 几个 sample (玩具 / 钥匙扣)\n• 现场 让 学生 摸 一摸!")

# ============================================================
# 22 · Live Demo 观察 任务
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔬 现场 看 3D 打印机!  Live Demo · Be a Scientist", MODERN)

# Top intro
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.80))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = MODERN; intro.line.width = Pt(2.5)
tb(s, 0.55, 1.05, 9.00, 0.35, "🧑‍🔬 你 现在 是 「小 科学家」 — 仔 细 观察 3D 打印机!",
   sz=14, b=True, c=MODERN, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.45, 9.00, 0.28, "Be a Tiny Scientist — observe the 3D printer carefully!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: 3D printer photo
photo_slot(s, 0.40, 1.95, 3.80, 3.05,
           "Bambu Lab / Prusa / 课堂 3D 打印机 照片",
           "Photo of classroom 3D printer in action",
           color=MODERN)

# RIGHT: 5 observation questions
panel(s, 4.30, 1.95, 5.30, 3.05, PURPLE)
panel_head(s, 4.30, 1.95, 5.30, PURPLE, "🔍 观察 任务 · 5 个 问题", sz=12)
obs_qs = [
    ("👁️", "它 在 用 什么 材料?", "What material?"),
    ("📚", "它 是 一层 一层 做 的 吗?", "Layer by layer?"),
    ("⏱️", "快 还是 慢?", "Fast or slow?"),
    ("👂", "声音 像 什么?", "What sound?"),
    ("🎨", "它 在 「画」 还是 「做」?", "Drawing or making?"),
]
for i, (em, cn, en) in enumerate(obs_qs):
    y = 2.50 + i * 0.52
    tb(s, 4.45, y, 0.45, 0.40, em, sz=18)
    tb(s, 4.95, y, 3.50, 0.32, cn, sz=11, b=True, c=DARK)
    tb(s, 4.95, y + 0.32, 3.50, 0.22, en, sz=8, c=GRAY)

# Higher-grade extension
ext = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.45))
ext.fill.solid(); ext.fill.fore_color.rgb = INK; ext.line.color.rgb = STAR; ext.line.width = Pt(2)
tb(s, 0.55, 5.16, 9.00, 0.32, "💎 高 年级: 「为什么 不能 一下 变 出 来?」",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "15 分钟 现场 Demo:\n\n准备 (老师 提前):\n• 3D 打印机 接上 电源 + 网络\n• 已经 加载 好 一个 简单 模型 (5-10 分钟 能 打 完 的 小 物 — 比如 钥匙 扣)\n• 投影 显示 屏幕 操作\n\n现场:\n• 启动 打印机\n• 让 学生 围观 (注意 安全 — 不 让 摸 高温 部件!)\n• 一边 打印 一边 问 5 个 观察 问题\n• 学生 在 booklet 上 记 答案\n\n候 打印 时:\n• 继续 讲解 / 玩 游戏\n• 让 学生 看 屏幕 上 的 「层数 计数」\n\n打印 完:\n• 让 学生 排队 摸 (等 冷 了)\n• 大家 一起 鼓掌\n\n安全 提示:\n• 喷头 200°C+ — 绝对 不 能 碰\n• 打印 中 不 要 让 头发 / 手 靠 近\n• 不 让 学生 自己 操 作 按钮")

# ============================================================
# 23 · 对比 总结 · 活字印刷 vs 3D 打印
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "⚖️ 对比 · 古代 vs 现代  Ancient vs Modern", PURPLE)
tb(s, 0.40, 0.85, 9.20, 0.30, "想 一 想 — 古代 印刷 和 3D 打印 一 样? 不 一 样?",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: 活字印刷
panel(s, 0.40, 1.25, 4.55, 3.20, ANCIENT)
panel_head(s, 0.40, 1.25, 4.55, ANCIENT, "📜 活字 印刷 (1000 年 前)", sz=13)
ancient_pts = [
    "🔤  印 「字」 — 在 纸 上",
    "🧱  字块 排好 + 刷墨 + 压 纸",
    "📚  平面 (2D)",
    "✋  手工 操作",
    "🇨🇳  毕昇 (中国 发明)",
]
for i, line in enumerate(ancient_pts):
    tb(s, 0.55, 1.85 + i * 0.50, 4.30, 0.42, line, sz=11, b=True, c=DARK)

# RIGHT: 3D 打印
panel(s, 5.05, 1.25, 4.55, 3.20, MODERN)
panel_head(s, 5.05, 1.25, 4.55, MODERN, "🖨️ 3D 打印 (现代!)", sz=13)
modern_pts = [
    "🎁  印 「东西」 — 真 物 体",
    "🍦  一层 一层 堆 材 料",
    "🧊  立体 (3D)",
    "🤖  机器 自动",
    "🌍  全 世 界 都 用",
]
for i, line in enumerate(modern_pts):
    tb(s, 5.20, 1.85 + i * 0.50, 4.30, 0.42, line, sz=11, b=True, c=DARK)

# Bottom — what's the same?
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.90))
band.fill.solid(); band.fill.fore_color.rgb = INK
band.line.color.rgb = STAR; band.line.width = Pt(2)
tb(s, 0.55, 4.72, 9.00, 0.28, "💡 一 样 在 哪 里?  Both have in common:",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.02, 9.00, 0.30, "📐 都 要 设计   ⚙️ 都 要 机器/工具   ✨ 都 把 「想法」 变 「真 实」!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.34, 9.00, 0.18, "Both: design first · use tools · turn ideas into reality!",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 比较 + 反思:\n• 老师 念 5 个 对比 点\n• 让 学生 找 「一样」 的 地方\n• 关键 insight: 「都 是 把 设计 变 真实」 — 这 就 是 「制造」!\n• 引出 项目: 你 今天 也 来 制造!")

# ============================================================
# 24 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3", "🎨 下午 3:00–4:30  ·  古代 + 未来 制造 工厂!", DAY, "🛠️")
n += 1; pn(s, n)

# Complete today's booklet — before project starts
s = booklet_slide(prs, day_num=2, day_topic_cn="活字 印刷 + 3D 打印", day_color=DAY)
n += 1; pn(s, n)

# ============================================================
# 25 · 2 Projects overview
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🛠️ 今天 的 项目 · 古代 + 未来!  2 Projects", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "上午 学 古代 + 下午 学 现代 — 现在 你 都 来 做 一 做!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

projects = [
    ("📜", "项目 1", "DIY 活字 印刷", "DIY Movable Type",
     "All", "个人/小组",
     "自己 刻 字 + 拓印 — 像 毕昇 一样!",
     "Carve a character + print it — like Bi Sheng!", ANCIENT),
    ("🎨", "项目 2", "3D 玩具 涂色 + 命名", "Color & Name 3D Toy",
     "All", "个人",
     "给 老师 3D 打印 的 玩具 涂色 + 取 名 + 介 绍",
     "Color the 3D printed toy + name it + present", MODERN),
]
card_w = 4.40; gap = 0.30
total = 2 * card_w + gap; start = (10 - total) / 2
for i, (em, num, cn_t, en_t, level, style, cn_d, en_d, cl) in enumerate(projects):
    x = start + i * (card_w + gap)
    panel(s, x, 1.30, card_w, 3.70, cl, lw=3)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.30), Inches(card_w), Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb = cl; head.line.fill.background()
    tb(s, x, 1.36, card_w, 0.40, num, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x, 1.95, card_w, 0.85, em, sz=60, a=PP_ALIGN.CENTER)
    tb(s, x, 2.90, card_w, 0.45, cn_t, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 3.35, card_w, 0.30, en_t, sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    # Level + style badges
    b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.30), Inches(3.75), Inches(1.0), Inches(0.32))
    b1.fill.solid(); b1.fill.fore_color.rgb = cl; b1.line.fill.background()
    tb(s, x+0.30, 3.78, 1.0, 0.26, level, sz=10, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+1.40), Inches(3.75), Inches(card_w-1.70), Inches(0.32))
    b2.fill.solid(); b2.fill.fore_color.rgb = WHITE; b2.line.color.rgb = cl; b2.line.width = Pt(1.5)
    tb(s, x+1.40, 3.78, card_w-1.70, 0.26, style, sz=9, b=True, c=cl, a=PP_ALIGN.CENTER)
    # Description
    tb(s, x+0.15, 4.20, card_w-0.30, 0.45, cn_d, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x+0.15, 4.65, card_w-0.30, 0.30, en_d, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Time tip
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.42))
tip.fill.solid(); tip.fill.fore_color.rgb = DAY; tip.line.fill.background()
tb(s, 0.55, 5.16, 9.00, 0.30, "⏱️ 30 min 项目 1 + 30 min 项目 2 + 30 min 分享!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "项目 安排:\n• 两 个 项目 都 做 (or 选 一个 重点 做)\n• 先 项目 1 (古代) → 然后 项目 2 (现代)\n• 分享 时 一起 比较\n\n材料 准备:\n• Project 1: foam sheets / carving stamp blocks / sponges / 墨 (or color pens)\n• Project 2: 提前 3D 打印 好 小 玩具 sample (每人 一个) + 涂色 笔")

# ============================================================
# 26 · Project 1 detail · DIY 活字印刷
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📜 项目 1 · DIY 活字 印刷  Carve & Print!", ANCIENT)

# LEFT: 4 steps
panel(s, 0.40, 0.95, 4.55, 4.15, ANCIENT)
panel_head(s, 0.40, 0.95, 4.55, ANCIENT, "📝 怎么 做  Steps", sz=12)
steps = [
    "1️⃣ 选 一个 字 (中/爱/福/家/AI...)",
    "2️⃣ 在 foam / 橡皮 上 画 字 (反 着!)",
    "3️⃣ 老师 帮 你 刻 出 来",
    "4️⃣ 上 墨 (或 颜料)",
    "5️⃣ 压 在 纸 上 — 印!",
    "6️⃣ 再 印 一 个 — 看 一样 吗?",
]
for i, line in enumerate(steps):
    tb(s, 0.55, 1.55 + i * 0.55, 4.30, 0.45, line, sz=11, b=True, c=DARK)

# RIGHT: word options + photo
panel(s, 5.05, 0.95, 4.55, 2.10, ORANGE)
panel_head(s, 5.05, 0.95, 4.55, ORANGE, "✏️ 选 一个 字 刻!  Pick a Character", sz=12)
chars = ["中", "爱", "福", "家", "AI"]
for i, ch in enumerate(chars):
    x = 5.20 + i * 0.85
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.55), Inches(0.75), Inches(1.30))
    card.fill.solid(); card.fill.fore_color.rgb = WARM
    card.line.color.rgb = ORANGE; card.line.width = Pt(2)
    tb(s, x, 1.75, 0.75, 0.90, ch, sz=30, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)

# RIGHT-bottom: materials
panel(s, 5.05, 3.20, 4.55, 1.90, PINK)
panel_head(s, 5.05, 3.20, 4.55, PINK, "🎨 材料  Materials", sz=12)
mats = [
    "🧽 Foam sheets / 橡皮 / 雕刻 块",
    "✂️ 刻 刀 (老师 帮!)",
    "🖌️ 墨 / 颜料",
    "📄 纸",
]
for i, line in enumerate(mats):
    tb(s, 5.20, 3.70 + i * 0.32, 4.30, 0.28, line, sz=10, b=True, c=DARK)

# Safety + tip bar
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.40), Inches(5.20), Inches(9.20), Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb = ANCIENT; tip.line.fill.background()
tb(s, 0.55, 5.26, 9.00, 0.30, "⚠️ 安全: 刻刀 老师 用! 学生 只 画 + 印  ·  你 是 「小 毕昇」!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "30 分钟:\n\n材料 选 一个:\n• Foam sheets (最 安全, 用 钝 头 工具 也 能 刻)\n• Carving stamp blocks (橡皮 砖 — 比较 软)\n• Speedy-cut 雕刻 块 (儿童 安全 版)\n• Sponge (剪 形 状, 不 用 刻 — 最 简单)\n\n步骤:\n1. 学生 画 字 (画 在 paper 上 翻 过 来 拓 到 foam)\n2. 老师 帮 刻 (或 让 学生 用 安全 工具)\n3. 上 墨 (印泥 或 水彩 颜料)\n4. 压 在 纸 上 — 揭开 看!\n\n讨论:\n• 「印 出 来 和 原来 画 的 一 样 吗?」 (反 的!)\n• 「为什么 字 要 反 着 刻?」 (因为 印 出来 才 是 正 的)\n\nK-2: 用 sponge 形状 + 颜料\n3-5: 真的 刻 字 / 自己 名字 首字母")

# ============================================================
# 27 · Project 2 detail · 3D 玩具 涂色
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎨 项目 2 · 3D 玩具 涂色 + 命名!  Color & Name", MODERN)

# LEFT: steps
panel(s, 0.40, 0.95, 4.55, 4.15, MODERN)
panel_head(s, 0.40, 0.95, 4.55, MODERN, "📝 怎么 做  Steps", sz=12)
steps2 = [
    "1️⃣ 选 一个 老师 3D 打印 的 玩具",
    "2️⃣ 想 一 想 — 它 是 谁? 叫 什么?",
    "3️⃣ 用 颜色 笔 涂 色",
    "4️⃣ 给 它 取 一个 名 字",
    "5️⃣ 写 / 说 它 的 故事",
    "6️⃣ 上 台 介绍 给 全 班!",
]
for i, line in enumerate(steps2):
    tb(s, 0.55, 1.55 + i * 0.55, 4.30, 0.45, line, sz=11, b=True, c=DARK)

# RIGHT: sample 3D-printed toys to choose from
panel(s, 5.05, 0.95, 4.55, 2.10, PURPLE)
panel_head(s, 5.05, 0.95, 4.55, PURPLE, "🎁 玩具 选 项  Choose Your Toy", sz=12)
toys = [
    ("🤖", "小 机器人"),
    ("🦖", "小 恐 龙"),
    ("🚀", "小 火 箭"),
    ("🐶", "小 动物"),
]
for i, (em, line) in enumerate(toys):
    col = i % 2; row = i // 2
    x = 5.20 + col * 2.20
    y = 1.55 + row * 0.70
    tb(s, x, y, 0.55, 0.55, em, sz=24)
    tb(s, x + 0.60, y + 0.10, 1.55, 0.42, line, sz=11, b=True, c=DARK)

# RIGHT-bottom: sentence frames
panel(s, 5.05, 3.20, 4.55, 1.90, GREEN)
panel_head(s, 5.05, 3.20, 4.55, GREEN, "🗣️ 介绍 句型  Tell us about it!", sz=12)
frames = [
    "「我 的 玩具 叫 ___」",
    "「它 是 3D 打印 的!」",
    "「它 是 一个 ___」",
    "「它 会 ___」  (高 年级)",
]
for i, line in enumerate(frames):
    tb(s, 5.20, 3.70 + i * 0.32, 4.30, 0.28, line, sz=10, b=True, c=DARK)

# Tip
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.40), Inches(5.20), Inches(9.20), Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb = MODERN; tip.line.fill.background()
tb(s, 0.55, 5.26, 9.00, 0.30, "✨ 没有 错 答 案 — 你 的 想象 力 是 最 棒 的!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "30 分钟:\n\n准备:\n• 老师 提前 3D 打印 一批 小 玩具 (每 人 一个)\n  - 小 机器人 / 恐龙 / 火箭 / 动物 / 钥匙 扣\n  - 建议 每个 大 约 2-3 小时 打印, 提前 几天 准备\n• 准备 涂色 笔 / 彩色 马克 笔 (注意: 要 能 涂 在 PLA 塑料 上!)\n  - Sharpie / Posca 笔 最 好\n\n活动 流程:\n• 5 min: 让 学生 选 玩具\n• 15 min: 涂色\n• 10 min: 想 名字 + 写 故事\n\nK-2: 重 涂色, 简单 句型\n3-5: 设计 玩具 用途 + 自己 写 故事")

# ============================================================
# 28 · Share + Goodbye
# ============================================================
s = share_close(prs, DAY,
    frames_cn=[
        "「我 的 玩具 叫 ___, 它 是 3D 打印 的!」",
        "「古代 印刷 和 3D 打印 都 ___」",
    ],
    frames_en="My toy is ___, it was 3D printed!  ·  Ancient + 3D printing both ___",
    next_day_cn="Day 3 · Machine Learning — 电脑 怎么 学?",
    next_day_en="Day 3 · How do computers learn?",
    next_emoji="🧠")
n += 1; pn(s, n)
notes(s, "15 分钟 分享:\n• 每个 学生 上 台 1 分钟\n• 用 句型 介绍 自己 的 玩具\n• 全班 比较: 古代 印刷 和 现代 3D 打印 — 一样? 不一样?\n• 老师 拍 照\n• 提醒 明天 主题: Machine Learning")

# ============================================================
out = os.path.join(os.path.dirname(__file__), "day2_3dprint.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
