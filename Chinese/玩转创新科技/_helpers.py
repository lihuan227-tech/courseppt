# Shared helpers for 玩转创新科技 unit (Day 1-5)
# Color palette + reusable layout primitives modeled on the 仰望星空 scripts.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ===== Unit palette =====
INK    = RGBColor(0x0F,0x1A,0x3A)
CYBER  = RGBColor(0x2A,0x47,0xE0)
NEON   = RGBColor(0x00,0xBC,0xD4)
GOLD   = RGBColor(0xFF,0xC1,0x07)
STAR   = RGBColor(0xF5,0xC2,0x42)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
CREAM  = RGBColor(0xFF,0xF8,0xE7)
WARM   = RGBColor(0xFF,0xF3,0xE0)
IMGBG  = RGBColor(0xE8,0xE8,0xF0)
RED    = RGBColor(0xC8,0x25,0x3E)
SKY    = RGBColor(0x42,0xA5,0xF5)

# Day accent palette (one per day)
AI_PURPLE     = RGBColor(0x6A,0x1B,0x9A)
PRINT_ORANGE  = RGBColor(0xFB,0x8C,0x00)
ML_GREEN      = RGBColor(0x2E,0x7D,0x32)
LIFE_TEAL     = RGBColor(0x00,0x79,0x6B)
FUTURE_PINK   = RGBColor(0xD8,0x1B,0x60)

# Secondary accents (used inside slides)
ORANGE = PRINT_ORANGE
PURPLE = AI_PURPLE
TEAL   = LIFE_TEAL
PINK   = FUTURE_PINK
GREEN  = ML_GREEN
EARTH  = RGBColor(0x1E,0x88,0xE5)


def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs


def ns(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def tb(s, l, t, w, h, txt, sz=18, b=False, c=DARK, a=None):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if a:
        p.alignment = a
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = 'KaiTi'
    return tf


def bg(s, c, prs):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    sp = sh._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)


def hb(s, txt, c=INK, t=0.15):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(t), Inches(9.4), Inches(0.55))
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    tb(s, 0.4, t+0.03, 9.2, 0.5, txt, sz=20, b=True, c=WHITE)


def pn(s, n):
    tb(s, 9.0, 5.05, 0.8, 0.30, str(n), sz=10, c=GRAY, a=PP_ALIGN.RIGHT)


def notes(s, text):
    nf = s.notes_slide.notes_text_frame
    lines = text.split("\n")
    nf.text = lines[0]
    for line in lines[1:]:
        p = nf.add_paragraph()
        p.text = line


def panel(s, l, t, w, h, color, fill=WHITE, lw=2.5):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.color.rgb = color
    p.line.width = Pt(lw)
    return p


def panel_head(s, l, t, w, color, txt, text_color=WHITE, sz=14):
    h = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.50))
    h.fill.solid()
    h.fill.fore_color.rgb = color
    h.line.fill.background()
    tb(s, l+0.15, t+0.07, w-0.3, 0.40, txt, sz=sz, b=True, c=text_color)


def div(prs, title, sub, color, emoji=""):
    """Section divider slide. Uses short title; subtitle is one line."""
    s = ns(prs)
    bg(s, color, prs)
    tb(s, 0.3, 1.50, 9.4, 1.0, f"{emoji} {title}", sz=44, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 0.3, 2.70, 9.4, 0.4, sub, sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    for x, y in [(0.8, 4.7), (1.8, 4.5), (7.8, 4.5), (8.6, 4.7)]:
        d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
        d.fill.solid()
        d.fill.fore_color.rgb = STAR
        d.line.fill.background()
    return s


def tianzi_box(s, x, y, size, char, color, pinyin=None, char_sz=130):
    """Tian-zi-ge (田字格) practice cell with character."""
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(3)
    mx = x + size/2
    my = y + size/2
    lw = 0.015
    v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx-lw/2), Inches(y), Inches(lw), Inches(size))
    v.fill.solid()
    v.fill.fore_color.rgb = LGRAY
    v.line.fill.background()
    h = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(my-lw/2), Inches(size), Inches(lw))
    h.fill.solid()
    h.fill.fore_color.rgb = LGRAY
    h.line.fill.background()
    tb(s, x, y, size, size, char, sz=char_sz, b=True, c=color, a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s, x, y+size+0.05, size, 0.30, pinyin, sz=12, b=True, c=GRAY, a=PP_ALIGN.CENTER)


def cover(prs, day_num, cn_title, en_title, emoji_row, day_color, inquiry_cn, inquiry_en):
    """Standard cover slide: unit name, day, title, emojis, inquiry question."""
    s = ns(prs)
    bg(s, INK, prs)
    tb(s, 0.5, 0.40, 9.0, 0.50, "🚀 玩转创新科技  Playing with Innovative Tech",
       sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 1.00, 9.0, 0.50, f"Day {day_num} · {cn_title}",
       sz=30, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 1.55, 9.0, 0.40, en_title,
       sz=16, c=LGRAY, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 2.30, 9.0, 0.70, emoji_row,
       sz=44, c=day_color, a=PP_ALIGN.CENTER)
    # Inquiry question card
    qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.50), Inches(7.6), Inches(1.55))
    qbox.fill.solid()
    qbox.fill.fore_color.rgb = day_color
    qbox.line.color.rgb = STAR
    qbox.line.width = Pt(3)
    tb(s, 1.3, 3.60, 7.4, 0.35, "🤔 今天的探究问题  Today's Inquiry",
       sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 1.3, 4.00, 7.4, 0.45, inquiry_cn, sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 1.3, 4.55, 7.4, 0.35, inquiry_en, sz=11, c=WARM, a=PP_ALIGN.CENTER)
    return s


def learning_goals(prs, header_color, goals):
    """Learning goals slide. goals is list of (number_emoji, cn_text, en_text, color)."""
    s = ns(prs)
    bg(s, CREAM, prs)
    hb(s, "🎯 今天的学习目标  Today's Learning Goals", header_color)
    tb(s, 0.4, 0.85, 9.2, 0.32, "上完这节课, 你会……  By the end, you'll be able to…",
       sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    for i, (num, cn, en, cl) in enumerate(goals):
        y = 1.30 + i * 0.95
        panel(s, 0.5, y, 9.0, 0.80, cl, fill=WHITE, lw=2.5)
        tb(s, 0.65, y+0.08, 0.6, 0.40, num, sz=20, b=True, c=cl)
        tb(s, 1.30, y+0.08, 8.0, 0.32, cn, sz=13, b=True, c=DARK)
        tb(s, 1.30, y+0.43, 8.0, 0.30, en, sz=10, c=GRAY)
    return s


def vocab_recognize(prs, header_color, big_emoji, cn, pinyin, en, example_cn, example_en, picture_hint):
    """我会认 slide — single word focus."""
    s = ns(prs)
    bg(s, CREAM, prs)
    hb(s, f"👀 我会认 · {cn}  I Can Read", header_color)
    # LEFT: big word card
    panel(s, 0.40, 0.95, 4.55, 3.30, header_color, fill=WARM)
    tb(s, 0.50, 1.10, 4.35, 0.80, big_emoji, sz=58, a=PP_ALIGN.CENTER)
    tb(s, 0.50, 1.95, 4.35, 0.70, cn, sz=44, b=True, c=header_color, a=PP_ALIGN.CENTER)
    tb(s, 0.50, 2.75, 4.35, 0.30, f"{pinyin}  ·  {en}", sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.50, 3.20, 4.35, 0.32, "👉 跟我读!  Read after me!", sz=12, b=True, c=header_color, a=PP_ALIGN.CENTER)
    # RIGHT: picture placeholder
    ib = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.10), Inches(0.95), Inches(4.50), Inches(3.30))
    ib.fill.solid()
    ib.fill.fore_color.rgb = IMGBG
    ib.line.color.rgb = header_color
    ib.line.width = Pt(2)
    tb(s, 5.10, 2.20, 4.50, 0.40, "📷 " + picture_hint, sz=13, c=LGRAY, a=PP_ALIGN.CENTER)
    tb(s, 5.10, 2.65, 4.50, 0.30, "图片位置 · Image placeholder", sz=10, c=LGRAY, a=PP_ALIGN.CENTER)
    # Bottom: example sentence
    ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.40), Inches(9.20), Inches(0.95))
    ex.fill.solid()
    ex.fill.fore_color.rgb = WHITE
    ex.line.color.rgb = header_color
    ex.line.width = Pt(2)
    tb(s, 0.55, 4.48, 9.00, 0.32, "📌 例句  Example", sz=11, b=True, c=header_color)
    tb(s, 0.55, 4.78, 9.00, 0.35, example_cn, sz=14, b=True, c=DARK)
    tb(s, 0.55, 5.13, 9.00, 0.28, example_en, sz=10, c=GRAY)
    return s


def vocab_write(prs, header_color, cn_phrase, en_word, chars):
    """我会写 slide. chars is list of (char, pinyin, stroke_count_note, mnemonic)."""
    s = ns(prs)
    bg(s, CREAM, prs)
    hb(s, f"✏️ 我会写 · {cn_phrase}  I Can Write · {en_word}", header_color)
    tb(s, 0.4, 0.85, 9.2, 0.32, f"一起来写「{cn_phrase}」!  Practice writing {cn_phrase} ({en_word})",
       sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    # Tian-zi-ge cells (1.10" each) on the left half
    n = len(chars)
    cell_size = 1.30
    gap_cells = 0.10
    total_w = n * cell_size + (n-1) * gap_cells
    start_x = 0.50
    for i, (ch, py, _stroke, _mn) in enumerate(chars):
        x = start_x + i * (cell_size + gap_cells)
        tianzi_box(s, x, 1.35, cell_size, ch, header_color, pinyin=py, char_sz=72)
    # Right side: instruction panel
    panel(s, 6.00, 1.35, 3.60, 3.30, header_color)
    panel_head(s, 6.00, 1.35, 3.60, header_color, "✏️ 怎么写  How to Write", sz=13)
    for i, (ch, py, stroke, mn) in enumerate(chars):
        y = 2.00 + i * 0.85
        tb(s, 6.15, y, 3.30, 0.30, f"📝 「{ch}」 — {py}", sz=12, b=True, c=DARK)
        tb(s, 6.15, y+0.30, 3.30, 0.50, f"{stroke}\n{mn}", sz=9, c=GRAY)
    # Bottom prompt (positioned above page number area)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.72), Inches(9.20), Inches(0.42))
    tip.fill.solid()
    tip.fill.fore_color.rgb = header_color
    tip.line.fill.background()
    tb(s, 0.55, 4.78, 9.00, 0.32, f"📓 在田字格里写 3 遍  ·  「我会写 {cn_phrase}!」",
       sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    return s


def share_close(prs, day_color, frames_cn, frames_en, next_day_cn, next_day_en, next_emoji="✨"):
    """Share + close slide."""
    s = ns(prs)
    bg(s, CREAM, prs)
    hb(s, "🎤 分享 + 再见!  Share + Goodbye", day_color)
    # Sentence frames panel
    fp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(2.05))
    fp.fill.solid()
    fp.fill.fore_color.rgb = INK
    fp.line.color.rgb = STAR
    fp.line.width = Pt(3)
    tb(s, 0.55, 1.05, 9.00, 0.35, "💬 今天我学了 — 句型  Today's Sentence Frames",
       sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
    for i, frame in enumerate(frames_cn):
        tb(s, 0.55, 1.50 + i*0.50, 9.00, 0.45, frame, sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 2.65, 9.00, 0.28, frames_en, sz=10, c=LGRAY, a=PP_ALIGN.CENTER)
    # Next day preview
    np = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.20), Inches(9.20), Inches(2.10))
    np.fill.solid()
    np.fill.fore_color.rgb = WHITE
    np.line.color.rgb = day_color
    np.line.width = Pt(2.5)
    tb(s, 0.55, 3.35, 9.00, 0.35, "🔮 明天见  See You Tomorrow", sz=14, b=True, c=day_color, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 3.80, 9.00, 0.50, f"{next_emoji}  {next_day_cn}", sz=18, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 4.30, 9.00, 0.30, next_day_en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 4.75, 9.00, 0.32, "👋 想一想: 你今天最喜欢的是什么?", sz=12, b=True, c=day_color, a=PP_ALIGN.CENTER)
    return s


# ===== Shared interactive helpers (used by Day 1, Day 3) =====

def photo_slot(s, l, t, w, h, suggestion_cn, suggestion_en, color=None):
    """Empty placeholder for a real photo (teacher will drop in later)."""
    if color is None:
        color = AI_PURPLE
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = IMGBG
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tb(s, l, t + 0.05, w, 0.25, "📷 真实图片 · Real Photo",
       sz=9, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    mid_y = t + h / 2 - 0.20
    tb(s, l, mid_y, w, 0.30, "在这里插入 →", sz=11, b=True, c=LGRAY, a=PP_ALIGN.CENTER)
    tb(s, l + 0.10, mid_y + 0.30, w - 0.20, 0.50, suggestion_cn,
       sz=10, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, l + 0.10, mid_y + 0.78, w - 0.20, 0.32, suggestion_en,
       sz=8, c=LGRAY, a=PP_ALIGN.CENTER)
    return box


def activity_box(s, l, t, w, h, prompt_cn, prompt_en, gesture_hint="", color=None):
    """Highlighted 「互动时间」 callout — flags when to engage the class.

    Adaptive layout based on h:
      • h <= 0.50 — compact: header + CN only (drop EN)
      • h <= 0.75 — header + CN + EN (no gesture line)
      • h >= 0.95 — header + CN + gesture + EN (full)
    """
    if color is None:
        color = AI_PURPLE
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = STAR
    box.line.width = Pt(2)

    # COMPACT mode: just header + Chinese prompt (fits in h ~0.55)
    if h <= 0.60:
        tb(s, l + 0.10, t + 0.04, w - 0.20, 0.22, "🎮 互动 · Activity",
           sz=10, b=True, c=STAR)
        tb(s, l + 0.15, t + 0.26, w - 0.30, 0.30, prompt_cn,
           sz=11, b=True, c=WHITE)
        return box

    # STANDARD mode: header + CN + EN (no gesture)
    # Content height: header 0.20 + cn 0.26 + en 0.22 = 0.68 → fits in h >= 0.75
    if h < 1.00:
        tb(s, l + 0.10, t + 0.03, w - 0.20, 0.20, "🎮 互动时间 · Activity",
           sz=10, b=True, c=STAR)
        tb(s, l + 0.15, t + 0.25, w - 0.30, 0.26, prompt_cn,
           sz=11, b=True, c=WHITE)
        tb(s, l + 0.15, t + 0.52, w - 0.30, 0.22, prompt_en,
           sz=8, c=WARM)
        return box

    # FULL mode (h >= 1.00): header + CN + gesture + EN
    # Content height: header 0.20 + cn 0.26 + gesture 0.22 + en 0.22 = 0.90 → fits in h >= 1.00
    tb(s, l + 0.10, t + 0.03, w - 0.20, 0.20, "🎮 互动时间 · Activity",
       sz=10, b=True, c=STAR)
    tb(s, l + 0.15, t + 0.25, w - 0.30, 0.26, prompt_cn,
       sz=11, b=True, c=WHITE)
    if gesture_hint:
        tb(s, l + 0.15, t + 0.52, w - 0.30, 0.22, gesture_hint,
           sz=9, b=True, c=STAR)
        tb(s, l + 0.15, t + 0.75, w - 0.30, 0.22, prompt_en,
           sz=8, c=WARM)
    else:
        tb(s, l + 0.15, t + 0.52, w - 0.30, 0.22, prompt_en,
           sz=8, c=WARM)
    return box
