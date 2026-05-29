#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
玩转 创新 科技 · Day 4: 从 古代 发明 到 现代 科技 — 科技 如何 改变 生活
3-session classroom deck for K-5 Chinese immersion summer camp.

探究 问题: 科技 怎么 改变 我们 的 生活?

我 会 认: 发明 / 科技 / 纸 / 过去 / 现在
我 会 写: 纸 / 过去 / 现在

Structure (matches Day 1 / Day 2 / Day 3):
  Session 1 (11:00–11:45) — 科技 如何 改变 生活? · Tech timeline + 造纸 术
  Session 2 (2:00–2:45)   — Bamboozle 复习 + 中文 词汇 + 过去 vs 现在
  Session 3 (3:00–4:30)   — Project: 造纸 体验 · Make Paper!
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = make_presentation()
DAY = LIFE_TEAL
ANCIENT = RGBColor(0xB8, 0x50, 0x42)  # terracotta — for ancient China sections
SEPIA   = RGBColor(0xA0, 0x6B, 0x3A)  # warm earth
n = 0


def arrow(s, x, y, w=0.30, h=0.30, color=DAY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ============================================================
# 1 · COVER
# ============================================================
cover(prs, 4, "Tech Changes Life", "科技 改变 生活 · 从 古代 到 现代",
      "🔥 🛞 📜 🖨 💡 🤖", DAY,
      "科技 怎么 改变 我们 的 生活?",
      "How does technology change our lives?")
n += 1; pn(prs.slides[-1], n)
notes(prs.slides[-1], "Day 4 · 一 天 主线:\n• 科技 解决 问题 → 科技 改变 生活 → 中国 造纸 → 自己 体验 古代 科技\n• Session 1: 科技 时间 旅行 + 中国 造纸 术\n• Session 2: Bamboozle 复习 + 中文 词汇 + 过去 vs 现在 配对\n• Session 3: 造纸 体验 项目")


# ============================================================
# 2 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1", "🌅 上午 11:00–11:45  ·  科技 如何 改变 生活?",
        DAY, "🔥"); n += 1; pn(s, n)


# ============================================================
# 3 · LEARNING GOALS
# ============================================================
s = learning_goals(prs, DAY, [
    ("1️⃣", "回 顾 这 周 学 过 的 科技 发明",
     "Review the inventions we learned this week", CYBER),
    ("2️⃣", "认识 中国 古代 发明 — 重 点 造纸 术",
     "Learn about ancient Chinese inventions — focus on paper", ANCIENT),
    ("3️⃣", "比 较 过去 vs 现在 — 生活 怎么 变 了",
     "Compare past vs present — how life changed", ORANGE),
    ("4️⃣", "懂 得 「科技 = 帮 人 解决 问题」",
     "Understand: tech is invented to solve human problems", PURPLE),
])
n += 1; pn(s, n)


# ============================================================
# 4 · HOOK — 科技 消失 的 一 天!
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.4, 0.45), (9.1, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.45, 9.4, 0.40, "⚠️ 假设 一下 ...",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.05), Inches(8.8), Inches(2.00))
tt.fill.solid(); tt.fill.fore_color.rgb = RED
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.25, 8.4, 0.85, "科技 消失 的 一 天!",
   sz=44, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.15, 8.4, 0.50, "The Day Tech Disappears!",
   sz=20, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.65, 8.4, 0.30, "今天 早 上 醒 来 — 所有 科技 都 不 见 了!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

tb(s, 0.3, 3.45, 9.4, 1.10, "🌅  😱  ❓  💭",
   sz=58, a=PP_ALIGN.CENTER)

tb(s, 0.3, 4.95, 9.4, 0.35, "🤔 你 会 怎么 办?",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "1 分钟 hype:\n• 用 戏 剧 性 的 语 气 说: 「今天 早 上 ...」\n• 让 学 生 想象 一下\n• 然后 切 到 下 一 张 看 没 有 什么 科技")


# ============================================================
# 5 · WHAT'S MISSING + TURN & TALK
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "❌ 没 有 这 些! · No More These!", RED)

tb(s, 0.4, 0.85, 9.2, 0.32, "想一想 — 没 有 这 些, 你 怎么 生活?",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "How would you live without these?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

missing = [
    ("💡", "电 灯", "Lights"),
    ("📱", "手 机", "Phones"),
    ("💻", "电 脑", "Computers"),
    ("📺", "电 视", "TV"),
    ("🚗", "汽 车", "Cars"),
    ("🌐", "Wi-Fi", "Internet"),
]
mw = 2.85; mgap_x = 0.20; mgap_y = 0.15; mh = 1.30
mstart_x = (10 - 3*mw - 2*mgap_x)/2
for i, (em, cn, en) in enumerate(missing):
    row = i // 3; col = i % 3
    x = mstart_x + col*(mw + mgap_x)
    y = 1.55 + row*(mh + mgap_y)
    panel(s, x, y, mw, mh, RED, fill=WHITE, lw=2.5)
    tb(s, x+0.05, y+0.05, 0.45, 0.40, "❌", sz=14, a=PP_ALIGN.LEFT)
    tb(s, x+0.55, y+0.10, 1.00, 0.85, em, sz=40, a=PP_ALIGN.LEFT)
    tb(s, x+1.65, y+0.20, mw-1.75, 0.45, cn, sz=18, b=True, c=RED)
    tb(s, x+1.65, y+0.65, mw-1.75, 0.30, en, sz=11, c=GRAY)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.45), Inches(9.20), Inches(1.00))
tt.fill.solid(); tt.fill.fore_color.rgb = DAY
tt.line.color.rgb = STAR; tt.line.width = Pt(2.5)
tb(s, 0.55, 4.52, 9.0, 0.32, "💬 Turn & Talk · 同 桌 讨论:",
   sz=13, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.85, 9.0, 0.30, "• 你 最 想念 什么?  • 什么 事 会 变 难?  • 你 还 能 怎么 生活?",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.18, 9.0, 0.22, "What do you miss most? What gets hard? How would you live?",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "5-8 分钟 互动:\n• 让 学 生 看 6 个 「失 去」 的 科技 — 引 起 共 鸣\n• Turn & Talk 3 分钟 — 同 桌 讨论\n• 老师 听 几 组 分 享\n• 总 结: 「科技 每 天 都 在 帮 助 我们 生活」\n• 引 出: 「可 是 很 久 以前 没 有 这 些 ... 那 怎么 办?」")


# ============================================================
# 6 · CONNECT — review Day 1/2/3
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔄 这 周 学 过 什么? · What We Learned This Week", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "我们 已经 认识 了 很 多 科技 — 一 起 回 顾!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

days = [
    ("Day 1", "🤖", "AI", "什么 是 AI · 做 聪明 AI 主人", AI_PURPLE),
    ("Day 2", "🖨️", "3D 打印", "古代 印刷 → 3D 打印", PRINT_ORANGE),
    ("Day 3", "🧠", "机器 学习", "机器 学习 + AI 训练 师", ML_GREEN),
]
dw = 2.95; dgap = 0.20
dtotal = 3*dw + 2*dgap; dstart = (10 - dtotal)/2
for i, (day, em, topic, sub, cl) in enumerate(days):
    x = dstart + i*(dw + dgap)
    panel(s, x, 1.30, dw, 3.40, cl, fill=WHITE, lw=3)
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.30), Inches(dw), Inches(0.55))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.40, dw, 0.40, day, sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x, 2.00, dw, 1.05, em, sz=64, a=PP_ALIGN.CENTER)
    tb(s, x, 3.10, dw, 0.50, topic, sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.70, dw-0.20, 0.65, sub, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

tb_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.90), Inches(9.20), Inches(0.55))
tb_bar.fill.solid(); tb_bar.fill.fore_color.rgb = DAY; tb_bar.line.fill.background()
tb(s, 0.55, 5.00, 9.0, 0.35, "💡 科技 一 直 在 进 步 — 今天 看 它 是 怎么 改变 生活 的!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 复 习:\n• 问 学 生: 「Day 1/2/3 学 了 什么?」\n• 让 学 生 大 声 喊 出 — AI / 3D / 机器 学 习\n• 总 结: 科技 一 直 在 进 步\n• 引 出 今天 主 题: 科技 怎么 一 步 一 步 改变 生活")


# ============================================================
# 7 · TECH TIMELINE — 科技 时间 旅行
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🚀 科技 时间 旅行 · Tech Time Travel", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "从 古代 到 现在 — 这 些 都 改变 了 生活!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

inventions = [
    ("🔥", "火",        ANCIENT),
    ("🛞", "轮 子",     ANCIENT),
    ("📜", "纸",        ANCIENT),
    ("🖨️", "印 刷",    SEPIA),
    ("☎️", "电 话",     CYBER),
    ("💡", "电 灯",     STAR),
    ("💻", "电 脑",     CYBER),
    ("🌐", "互 联 网",   DAY),
    ("🤖", "AI",        AI_PURPLE),
]
iw = 2.65; igap_x = 0.20; igap_y = 0.18; ih = 1.10
istart_x = (10 - 3*iw - 2*igap_x)/2
for i, (em, cn, cl) in enumerate(inventions):
    row = i // 3; col = i % 3
    x = istart_x + col*(iw + igap_x)
    y = 1.30 + row*(ih + igap_y)
    panel(s, x, y, iw, ih, cl, fill=WHITE, lw=2.5)
    tb(s, x+0.10, y+0.18, 0.85, 0.75, em, sz=42, a=PP_ALIGN.LEFT)
    tb(s, x+1.00, y+0.30, iw-1.10, 0.55, cn, sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
    nb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+iw-0.45), Inches(y+0.10), Inches(0.35), Inches(0.35))
    nb.fill.solid(); nb.fill.fore_color.rgb = cl; nb.line.fill.background()
    tb(s, x+iw-0.45, y+0.13, 0.35, 0.30, str(i+1), sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)

tb(s, 0.4, 5.10, 9.2, 0.32, "🌅 古代  →  →  →  →  现代 🌃",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 浏览:\n• 让 学 生 大 声 念 每 个 发明\n• 问: 还 知 道 其他 发明 吗?\n• 强 调: 「从 古代 到 现代 — 一 步 一 步」\n• 引 出: 一 个 一 个 看 详 细!")


# ============================================================
# INVENTION DEEP-DIVE — 9 slides, one per invention
# Each slide: emoji + name, WHEN/WHERE/WHO info, kid story, problem→solution
# ============================================================
def invention_slide(emoji, cn_name, en_name, when, where_cn, where_flag,
                    who_cn, who_en, story_cn, story_en, problem, change, color):
    s = ns(prs); bg(s, CREAM, prs)
    hb(s, f"{emoji} {cn_name} · {en_name}", color)

    # TOP ROW: 3 info cards (When · Where · Who)
    info_items = [
        ("📅", "什么 时候?", "When?", when, None),
        ("🌍", "在 哪里?",  "Where?", f"{where_flag} {where_cn}", None),
        ("👤", "谁?",       "Who?",   who_cn, who_en),
    ]
    iw = 2.95; igap = 0.20
    istart = (10 - 3*iw - 2*igap)/2
    for i, (em, cn_label, en_label, value, value_en) in enumerate(info_items):
        x = istart + i*(iw + igap)
        panel(s, x, 0.95, iw, 1.55, color, fill=WHITE, lw=2.5)
        # Label header
        tb(s, x, 1.05, iw, 0.40, em, sz=24, a=PP_ALIGN.CENTER)
        tb(s, x, 1.45, iw, 0.30, cn_label, sz=12, b=True, c=color, a=PP_ALIGN.CENTER)
        tb(s, x+0.05, 1.75, iw-0.10, 0.22, en_label, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
        # Value
        tb(s, x+0.10, 2.00, iw-0.20, 0.30, value,
           sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
        if value_en:
            tb(s, x+0.10, 2.28, iw-0.20, 0.20, value_en,
               sz=8, c=GRAY, a=PP_ALIGN.CENTER)

    # MIDDLE: story panel (full width)
    story_panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.40), Inches(2.70), Inches(9.20), Inches(1.50))
    story_panel.fill.solid(); story_panel.fill.fore_color.rgb = WARM
    story_panel.line.color.rgb = color; story_panel.line.width = Pt(2.5)
    tb(s, 0.55, 2.78, 9.0, 0.30, "📖 小 故 事 · The Story",
       sz=12, b=True, c=color, a=PP_ALIGN.LEFT)
    tb(s, 0.55, 3.10, 9.0, 0.50, story_cn,
       sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, 0.55, 3.70, 9.0, 0.45, story_en,
       sz=10, c=GRAY, a=PP_ALIGN.LEFT)

    # BOTTOM: problem → change
    prob_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.40), Inches(4.40), Inches(4.30), Inches(1.00))
    prob_box.fill.solid(); prob_box.fill.fore_color.rgb = RED
    prob_box.line.fill.background()
    tb(s, 0.55, 4.48, 4.10, 0.28, "😰 解决 的 问题:",
       sz=11, b=True, c=STAR, a=PP_ALIGN.LEFT)
    tb(s, 0.55, 4.78, 4.10, 0.55, problem,
       sz=13, b=True, c=WHITE, a=PP_ALIGN.LEFT)

    # Arrow
    arrow(s, 4.75, 4.78, w=0.50, h=0.40, color=color)

    change_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(5.30), Inches(4.40), Inches(4.30), Inches(1.00))
    change_box.fill.solid(); change_box.fill.fore_color.rgb = color
    change_box.line.fill.background()
    tb(s, 5.45, 4.48, 4.10, 0.28, "✨ 改变 了 什么:",
       sz=11, b=True, c=WARM, a=PP_ALIGN.LEFT)
    tb(s, 5.45, 4.78, 4.10, 0.55, change,
       sz=13, b=True, c=WHITE, a=PP_ALIGN.LEFT)

    return s


# Invention 1 — 🔥 火
s = invention_slide(
    "🔥", "火", "Fire",
    "大约 150 万 年 前",
    "远古 人类", "🌍",
    "远古 人类", "Early Humans",
    "雷电 点 燃 了 树 — 早期 人类 学 会 保留 火 种, 后来 学 会 摩擦 木 头 生火!",
    "Lightning lit trees — early humans kept the flames, then learned to make fire by rubbing wood.",
    "太 冷 + 生 食 难 吃",
    "可以 煮 饭 / 取 暖 / 吓 走 野 兽",
    ANCIENT,
)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 火 是 最 早 + 最 重要 的 发现\n• 不 是 「发明」 — 是 学 会 用\n• 问 学生: 你 觉得 古 人 怎么 「保留」 火 种?")


# Invention 2 — 🛞 轮子
s = invention_slide(
    "🛞", "轮 子", "Wheel",
    "约 公元 前 3500 年",
    "美索 不达 米亚", "🇮🇶",
    "苏 美 尔 人", "Sumerians",
    "最 早 的 轮子 不 是 用 来 搬 东 西 — 是 用 来 做 陶 器! 后来 才 装 到 车 上。",
    "The first wheel wasn't for carts — it was for making pottery! Later people put it on carts.",
    "搬 重 物 太 累",
    "可以 快速 搬 运 + 出 行",
    ANCIENT,
)
n += 1; pn(s, n)
notes(s, "2 分钟 有 趣 事 实:\n• 轮子 最 早 不 是 车轮 — 是 陶 器 转 盘!\n• 这 让 学 生 觉 得 「啊?!」")


# Invention 3 — 📜 纸
s = invention_slide(
    "📜", "纸", "Paper",
    "公元 105 年",
    "中国", "🇨🇳",
    "蔡 伦", "Cài Lún",
    "蔡 伦 用 树 皮 + 旧 布 + 旧 鱼 网 做 纸 — 又 轻 又 便 宜!",
    "Cai Lun used tree bark + old cloth + old fishing nets to make paper — light and cheap!",
    "写 字 工具 太 重 / 太 贵",
    "知识 可以 保 存 + 传 播",
    ANCIENT,
)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 重 点 — 这 是 中国 「四 大 发明」 之 一\n• 强 调: 蔡 伦 「改 进」 纸 (不 是 完 全 发 明)\n• 这 张 跟 后面 的 造纸 项目 直接 连 接")


# Invention 4 — 🖨️ 印刷
s = invention_slide(
    "🖨️", "印 刷", "Printing",
    "公元 1040 年",
    "中国", "🇨🇳",
    "毕 昇", "Bì Shēng",
    "他 用 泥土 做 「活 字」 — 一 个 一 个 字, 可以 反 复 使 用 印 很 多 本 书!",
    "He made movable type from clay — single characters that could be reused to print many books!",
    "一 本 一 本 手 抄 太 慢",
    "书 变 多 — 知识 传 得 快",
    SEPIA,
)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 中国 「四 大 发明」 之 二\n• 提 醒: 这 跟 Day 2 学 过 的 一 样!\n• 让 学 生 大 声 喊: 「毕 昇!」")


# Invention 5 — ☎️ 电话
s = invention_slide(
    "☎️", "电 话", "Telephone",
    "1876 年",
    "美 国", "🇺🇸",
    "贝 尔", "A. G. Bell",
    "贝 尔 发明 电话 后, 第 一 句 话 是: 「Watson, come here, I want to see you!」",
    "After inventing the phone, Bell's first words were: 'Watson, come here, I want to see you!'",
    "太 远 不 能 说 话",
    "可以 跟 远 方 的 人 立 刻 说 话",
    CYBER,
)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 第 一 句 话 故 事 — 学 生 会 觉 得 有 趣\n• 问: 没 有 电话 你 怎么 跟 远 方 的 朋友 说话?")


# Invention 6 — 💡 电灯
s = invention_slide(
    "💡", "电 灯", "Light Bulb",
    "1879 年",
    "美 国", "🇺🇸",
    "爱 迪 生", "T. Edison",
    "爱迪生 试 了 1000 多 种 材料 — 才 找 到 能 用 的 灯 丝!",
    "Edison tested over 1,000 materials before finding one that worked as a filament!",
    "晚 上 太 黑",
    "晚 上 也 能 学 习 / 工 作",
    STAR,
)
n += 1; pn(s, n)
notes(s, "2 分钟 重 要 教 学:\n• 1000 多 次 失 败 — 但 没 放 弃\n• 教 学 点: 「失 败 是 学 习 的 一 部 分」\n• 联 想: AI 也 是 不 断 改 进 (Day 3 主 题)")


# Invention 7 — 💻 电脑
s = invention_slide(
    "💻", "电 脑", "Computer",
    "1946 年",
    "美 国", "🇺🇸",
    "ENIAC 团 队", "ENIAC Team",
    "第 一 台 电 脑 叫 ENIAC — 有 一 个 房 间 那 么 大, 重 30 吨!",
    "The first computer (ENIAC) was as big as a room and weighed 30 tons!",
    "算 复 杂 数 学 太 慢",
    "算 得 超 快 + 能 做 很 多 事",
    CYBER,
)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 让 学 生 想象: 一 台 电脑 = 你 们 的 教 室 那 么 大!\n• 30 吨 = 比 大 象 还 重\n• 对 比 现在 的 手 机 — 多 小 多 快!")


# Invention 8 — 🌐 互联网
s = invention_slide(
    "🌐", "互 联 网", "Internet",
    "1969 年",
    "美 国", "🇺🇸",
    "ARPANET 团 队", "ARPANET Team",
    "第 一 条 网络 信息 本 来 是 「LOGIN」 — 但 打 到 「LO」 电 脑 就 死 机 了!",
    "The first internet message was supposed to be 'LOGIN' — but it crashed after just 'LO'!",
    "信息 传 得 太 慢",
    "全 世 界 即 时 连 接",
    DAY,
)
n += 1; pn(s, n)
notes(s, "2 分钟 故 事:\n• 第 一 条 网络 信息 只 有 「LO」!\n• 让 学 生 笑\n• 引 导: 现在 互联网 一 秒 钟 传 几 亿 条 信息")


# Invention 9 — 🤖 AI
s = invention_slide(
    "🤖", "AI", "Artificial Intelligence",
    "1956 年",
    "美 国", "🇺🇸",
    "麦 卡 锡", "John McCarthy",
    "一 群 科学 家 开 会 讨论 「机器 能 不 能 像 人 一 样 思考」 — 那 次 第 一 次 用 了 「AI」 这 个 词!",
    "Scientists held a meeting asking 'Can machines think like humans?' — that's when 'AI' was first named!",
    "工 作 太 复 杂",
    "电脑 可以 「学 习」 + 帮 我 们 做 事",
    AI_PURPLE,
)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 联 系 Day 1 + Day 3 — AI / 机器 学习\n• 「AI」 这 个 词 是 1956 年 才 有 的!\n• 提 问: 你 现在 用 过 哪 些 AI? (复 习 Day 3)")


# ============================================================
# 8 · PROBLEMS TECH SOLVES — 科技 解决 问题
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔧 科技 = 解决 问题 · Tech Solves Problems", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "每 个 发明 都 解决 了 一 个 问题!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "Every invention solves a real problem",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

problems = [
    ("🔥", "火",       "😰 太 冷 / 生 食",     ANCIENT),
    ("🛞", "轮 子",    "😰 搬 东西 太 累",     ANCIENT),
    ("📜", "纸",       "😰 不 方便 记 知识",   ANCIENT),
    ("☎️", "电 话",    "😰 太 远 不 能 说",    CYBER),
    ("🌐", "互联网",   "😰 信息 传 得 慢",     DAY),
    ("🤖", "AI",       "😰 工作 太 复杂",     AI_PURPLE),
]
pw = 2.95; pgap_x = 0.18; pgap_y = 0.20; ph = 1.50
pstart_x = (10 - 3*pw - 2*pgap_x)/2
for i, (em, tech, problem, cl) in enumerate(problems):
    row = i // 3; col = i % 3
    x = pstart_x + col*(pw + pgap_x)
    y = 1.55 + row*(ph + pgap_y)
    panel(s, x, y, pw, ph, cl, fill=WHITE, lw=2.5)
    tb(s, x+0.10, y+0.12, 0.75, 0.50, em, sz=28, a=PP_ALIGN.LEFT)
    tb(s, x+0.85, y+0.18, pw-0.95, 0.40, tech, sz=17, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x, y+0.65, pw, 0.25, "↑ 解决 ↑", sz=10, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, y+0.95, pw-0.20, 0.45, problem, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
tk.fill.solid(); tk.fill.fore_color.rgb = DAY
tk.line.color.rgb = STAR; tk.line.width = Pt(2.5)
tb(s, 0.55, 5.10, 9.0, 0.30, "💡 科技 = 帮 人 解决 问题!  Tech = solving human problems!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5-7 分钟 讨论:\n• 让 学 生 想: 没 有 这 个 科技, 我 们 会 遇 到 什么 困 难?\n• 强 调: 每 个 发明 都 在 「解决 真 实 的 问题」")


# ============================================================
# 9 · ANCIENT WRITING — 以 前 的 人 写 字
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📜 以前 的 人 写 字 · How Ancient People Wrote", ANCIENT)

tb(s, 0.4, 0.85, 9.2, 0.30, "在 纸 没 有 以前 — 古 人 写 在 什么 上?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

surfaces = [
    ("🪨", "石 头", "Stone",
     "刻 字 — 太 重!\nCarve · too heavy!"),
    ("🎋", "竹 简", "Bamboo Strips",
     "穿 绳 — 太 重!\nBundles · still heavy!"),
    ("🧵", "丝 绸", "Silk",
     "写 字 — 太 贵!\nWrite on it · too expensive!"),
]
sw = 2.85; sgap = 0.25
stotal = 3*sw + 2*sgap; sstart = (10 - stotal)/2
for i, (em, cn, en, problem) in enumerate(surfaces):
    x = sstart + i*(sw + sgap)
    panel(s, x, 1.40, sw, 2.85, ANCIENT, fill=WHITE, lw=3)
    tb(s, x, 1.55, sw, 0.85, em, sz=56, a=PP_ALIGN.CENTER)
    tb(s, x, 2.50, sw, 0.45, cn, sz=20, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
    tb(s, x, 2.95, sw, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.15), Inches(3.40), Inches(sw-0.30), Inches(0.70))
    pb.fill.solid(); pb.fill.fore_color.rgb = WARM; pb.line.fill.background()
    bx = s.shapes.add_textbox(Inches(x+0.15), Inches(3.45), Inches(sw-0.30), Inches(0.65))
    tf = bx.text_frame; tf.word_wrap = True
    lines = problem.split("\n")
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = lines[0]; r0.font.size = Pt(11); r0.font.bold = True
    r0.font.color.rgb = ANCIENT; r0.font.name = 'KaiTi'
    for line in lines[1:]:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line; r.font.size = Pt(8); r.font.color.rgb = GRAY; r.font.name = 'KaiTi'

pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.50), Inches(9.20), Inches(0.95))
pb.fill.solid(); pb.fill.fore_color.rgb = ANCIENT
pb.line.color.rgb = STAR; pb.line.width = Pt(2.5)
tb(s, 0.55, 4.60, 9.0, 0.40, "😰 问题: 太 重 · 太 贵 · 不 方便!",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.05, 9.0, 0.30, "Problem: too heavy · too expensive · not convenient",
   sz=11, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟:\n• 让 学 生 想象: 你 的 课本 是 石头 做 的 — 怎么 带?\n• 竹 简 — 一 本 书 要 一 大 堆 竹片!\n• 丝 绸 — 又 贵 又 滑\n• 引 出 下 一 张: 这 时 候 — 蔡 伦 来 了!")


# ============================================================
# 10 · 蔡伦 改进 造纸 术
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🏛️ 蔡 伦 改 进 造 纸 术 · Cai Lun Improved Paper", ANCIENT)

tb(s, 0.4, 0.85, 9.2, 0.30, "一 个 中国 古人 — 改变 了 全 世界!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "One ancient Chinese man — changed the whole world!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

panel(s, 0.40, 1.55, 4.55, 3.55, ANCIENT, fill=WHITE, lw=3)
tb(s, 0.40, 1.70, 4.55, 1.30, "👨‍🔬",
   sz=110, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.05, 4.55, 0.55, "蔡 伦",
   sz=30, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.60, 4.55, 0.32, "Cài Lún · ~ 105 AD",
   sz=12, b=True, c=GRAY, a=PP_ALIGN.CENTER)
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.40), Inches(4.05), Inches(2.55), Inches(0.85))
cb.fill.solid(); cb.fill.fore_color.rgb = STAR; cb.line.fill.background()
tb(s, 1.50, 4.10, 2.35, 0.32, "🇨🇳 中国 古代",
   sz=13, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 1.50, 4.45, 2.35, 0.32, "Ancient China",
   sz=10, c=INK, a=PP_ALIGN.CENTER)

panel(s, 5.15, 1.55, 4.45, 3.55, DAY, fill=WHITE, lw=3)
panel_head(s, 5.15, 1.55, 4.45, DAY, "💡 他 做 了 什么?  What did he do?", sz=12)

cl_facts = [
    ("📜", "改 进 了 造 纸 术"),
    ("🌱", "用 树皮 + 旧 布 + 麻"),
    ("🪶", "纸 又 轻 + 又 便 宜"),
    ("📚", "更 多 人 可以 写 字 + 读 书!"),
]
for i, (em, txt) in enumerate(cl_facts):
    y = 2.25 + i*0.65
    tb(s, 5.30, y, 0.50, 0.50, em, sz=24, a=PP_ALIGN.LEFT)
    tb(s, 5.85, y+0.08, 3.65, 0.40, txt, sz=13, b=True, c=DARK)

emp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.20), Inches(9.20), Inches(0.30))
emp.fill.solid(); emp.fill.fore_color.rgb = ANCIENT; emp.line.fill.background()
tb(s, 0.55, 5.22, 9.0, 0.25, "🌟 造 纸 术 = 中国 「四 大 发明」 之 一!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 故 事:\n• 蔡 伦 — 东汉 时代 的 一 位 官 员\n• 不 是 他 「发 明」 了 纸 — 是 他 「改 进」 — 让 纸 又 便 宜 又 好 用\n• 强 调: 这 是 中国 「四 大 发明」 之 一 (其 他: 火 药, 指 南 针, 印刷 术)")


# ============================================================
# 11 · 有 了 纸 以 后 ...
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "✨ 有 了 纸 以 后 · After Paper Came", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "纸 让 全 世界 都 变 了!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

benefits = [
    ("✏️", "可以 写 字", "Can write easily", CYBER),
    ("📚", "可以 做 书", "Can make books", ORANGE),
    ("🧠", "可以 保 存 知识", "Can save knowledge", PURPLE),
    ("👨‍🎓", "更 多 人 学 习", "More people can learn", DAY),
]
bw = 4.40; bgap_x = 0.20; bgap_y = 0.20; bh = 1.45
bstart_x = (10 - 2*bw - bgap_x)/2
for i, (em, cn, en, cl) in enumerate(benefits):
    row = i // 2; col = i % 2
    x = bstart_x + col*(bw + bgap_x)
    y = 1.40 + row*(bh + bgap_y)
    panel(s, x, y, bw, bh, cl, fill=WHITE, lw=3)
    tb(s, x+0.20, y+0.20, 1.00, 1.00, em, sz=44, a=PP_ALIGN.LEFT)
    tb(s, x+1.35, y+0.25, bw-1.50, 0.50, cn, sz=18, b=True, c=cl)
    tb(s, x+1.35, y+0.80, bw-1.50, 0.32, en, sz=11, c=GRAY)

qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.80))
qb.fill.solid(); qb.fill.fore_color.rgb = ANCIENT
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 4.75, 9.0, 0.40, "🤔 如果 没 有 纸 — 学校 会 怎么 样?",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.18, 9.0, 0.22, "What would school be like without paper?",
   sz=10, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 讨论:\n• 重 点: 纸 = 知识 传 播 的 工具\n• 让 学 生 想象: 没 纸 → 没 课本, 没 作业 本, 没 故事 书\n• 桥 到 下 一 张: 「纸 + 印刷 = ?」")


# ============================================================
# 12 · QUICK CONNECTION — 纸 + 印刷
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📖 纸 + 印刷 = ? · Paper + Print = ?", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "Day 2 我们 学 过 印刷 — 加 上 纸, 会 发 生 什么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

panel(s, 0.40, 1.40, 2.50, 2.35, ANCIENT, fill=WHITE, lw=3)
tb(s, 0.40, 1.55, 2.50, 1.05, "📜", sz=70, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.65, 2.50, 0.45, "纸",
   sz=24, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.20, 2.50, 0.30, "Paper", sz=11, c=GRAY, a=PP_ALIGN.CENTER)

tb(s, 2.90, 2.30, 0.50, 0.80, "+", sz=44, b=True, c=DARK, a=PP_ALIGN.CENTER)

panel(s, 3.45, 1.40, 2.50, 2.35, PRINT_ORANGE, fill=WHITE, lw=3)
tb(s, 3.45, 1.55, 2.50, 1.05, "🖨️", sz=70, a=PP_ALIGN.CENTER)
tb(s, 3.45, 2.65, 2.50, 0.45, "印 刷",
   sz=24, b=True, c=PRINT_ORANGE, a=PP_ALIGN.CENTER)
tb(s, 3.45, 3.20, 2.50, 0.30, "Print", sz=11, c=GRAY, a=PP_ALIGN.CENTER)

tb(s, 5.95, 2.30, 0.50, 0.80, "=", sz=44, b=True, c=DARK, a=PP_ALIGN.CENTER)

panel(s, 6.50, 1.40, 3.10, 2.35, DAY, fill=WHITE, lw=3)
tb(s, 6.50, 1.55, 3.10, 1.05, "🌍📚", sz=58, a=PP_ALIGN.CENTER)
tb(s, 6.50, 2.65, 3.10, 0.45, "知识 传 播!",
   sz=20, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 6.50, 3.20, 3.10, 0.30, "Knowledge spreads!", sz=10, c=GRAY, a=PP_ALIGN.CENTER)

results = [
    ("📚", "更 多 书"),
    ("👨‍🎓", "更 多 学 习"),
    ("⚡", "知识 传 得 快"),
]
rw = 2.90; rgap = 0.20
rtotal = 3*rw + 2*rgap; rstart = (10 - rtotal)/2
for i, (em, txt) in enumerate(results):
    x = rstart + i*(rw + rgap)
    panel(s, x, 3.95, rw, 1.05, DAY, fill=WARM, lw=2)
    tb(s, x+0.15, 4.05, 0.80, 0.85, em, sz=32, a=PP_ALIGN.LEFT)
    tb(s, x+1.00, 4.30, rw-1.10, 0.45, txt, sz=14, b=True, c=DAY)

tb(s, 0.4, 5.10, 9.2, 0.30, "💡 一 个 发明 + 一 个 发明 = 更 大 的 变化!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3-5 分钟 桥 接:\n• 简 单 连 接 — 不 重 复 Day 2 内 容\n• 重 点: 纸 + 印刷 = 知识 传 播 速 度 飞 跃")


# ============================================================
# 13 · DISCUSSION — 哪 个 发明 最 改变 生活?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🗳️ 你 投 票! · You Vote!", DAY)

tb(s, 0.4, 0.85, 9.2, 0.32, "哪 个 发明 最 改变 生活? 为 什么?",
   sz=15, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.20, 9.2, 0.26, "Which invention changed life the most? Why?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

votes = [
    ("🔥", "火", ANCIENT),
    ("📜", "纸", ANCIENT),
    ("💡", "电 灯", STAR),
    ("☎️", "电 话", CYBER),
    ("💻", "电 脑", CYBER),
    ("🤖", "AI", AI_PURPLE),
]
vw = 2.85; vgap_x = 0.20; vgap_y = 0.15; vh = 1.25
vstart_x = (10 - 3*vw - 2*vgap_x)/2
for i, (em, cn, cl) in enumerate(votes):
    row = i // 3; col = i % 3
    x = vstart_x + col*(vw + vgap_x)
    y = 1.65 + row*(vh + vgap_y)
    panel(s, x, y, vw, vh, cl, fill=WHITE, lw=2.5)
    tb(s, x+0.15, y+0.20, 0.85, 0.85, em, sz=42, a=PP_ALIGN.LEFT)
    tb(s, x+1.10, y+0.35, vw-1.20, 0.55, cn, sz=20, b=True, c=cl)
    tb(s, x+vw-0.55, y+0.40, 0.40, 0.40, "☐", sz=22, c=cl, a=PP_ALIGN.CENTER)

mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.90))
mb.fill.solid(); mb.fill.fore_color.rgb = DAY
mb.line.color.rgb = STAR; mb.line.width = Pt(3)
tb(s, 0.55, 4.65, 9.0, 0.35, "💡 科技 不 是 突 然 出 现 — 是 一 步 一 步 改进 出 来 的!",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.05, 9.0, 0.30, "Tech doesn't appear suddenly — it improves step by step!",
   sz=10, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 讨论:\n• 让 学 生 举 手 投 票\n• 鼓 励 说 出 「为 什么」\n• 老师 总 结: 科技 是 「人 为 了 让 生活 更 方便 / 安全 / 快」 而 发明 的")


# ============================================================
# 14 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2", "📖 下午 2:00–2:45  ·  Bamboozle 复习 + 中文 词汇",
        DAY, "📚"); n += 1; pn(s, n)


# ============================================================
# 15 · BAMBOOZLE REVIEW (placeholder)
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.5, 0.5), (9.0, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.55, 9.4, 0.45, "🎮 GAME TIME!",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.30), Inches(8.8), Inches(2.30))
tt.fill.solid(); tt.fill.fore_color.rgb = DAY
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.55, 8.4, 0.85, "Bamboozle 复 习!",
   sz=44, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.50, 8.4, 0.45, "Review Game",
   sz=20, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 3.00, 8.4, 0.40, "复 习 今天 上午 学 的 知识!",
   sz=14, c=WARM, a=PP_ALIGN.CENTER)

tb(s, 0.3, 3.95, 9.4, 0.95, "🔥  📜  ✏️  📚  🌍  🚀",
   sz=52, a=PP_ALIGN.CENTER)

tb(s, 0.3, 5.00, 9.4, 0.40, "💡 老师 请 打 开 Bamboozle 游戏 链接",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "15-20 分钟 Bamboozle 复 习:\n• 老师 提前 准 备 好 Bamboozle 链 接\n• 复 习 题目 (建 议):\n  - 为 什么 火 重要?\n  - 蔡 伦 改 进 了 什么?\n  - 古 人 写 在 什么 上?\n  - 纸 解决 了 什么 问题?\n  - 没 有 纸 学校 会 怎 样?\n  - 科技 = 帮 人 做 什么?\n• 重 点 复 习 内容 理 解, 不 只 是 字 词")


# ============================================================
# 16-20 · 我 会 认 (vocabulary recognition) — 5 words
# ============================================================
recognize_words = [
    ("💡", "发明", "fā míng", "Invention",
     "蔡 伦 发 明 了 造 纸 术。", "Cai Lun invented paper making.",
     "灯 泡 / 创 新 / idea", PURPLE),
    ("🚀", "科技", "kē jì", "Technology",
     "科技 改变 了 我们 的 生活。", "Tech changed our lives.",
     "机器 / 电脑 / 火箭", CYBER),
    ("📜", "纸", "zhǐ", "Paper",
     "我 用 纸 写 字。", "I use paper to write.",
     "白 纸 / 书 / 卷 轴", ANCIENT),
    ("⏰", "过去", "guò qù", "Past",
     "过去 没有 电脑。", "In the past, there were no computers.",
     "古 代 / 老 照片 / 时 钟 倒 转", ORANGE),
    ("📱", "现在", "xiàn zài", "Now",
     "现在 我们 用 电脑。", "Now we use computers.",
     "今 天 / 城市 / 手机", DAY),
]
for em, cn, py, en, ex_cn, ex_en, hint, cl in recognize_words:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)


# ============================================================
# 21-23 · 我 会 写 (writing practice) — 纸 / 过去 / 现在
# ============================================================
s = vocab_write(prs, ANCIENT, "纸", "Paper",
                [("纸", "zhǐ", "7 笔", "「纟」 + 「氏」 — 古代 用 丝 做 纸")])
n += 1; pn(s, n)

s = vocab_write(prs, ORANGE, "过去", "Past",
                [("过", "guò", "6 笔", "「辶」 + 「寸」 — 走 过 去"),
                 ("去", "qù",  "5 笔", "上 「土」 下 「厶」 — 离 开")])
n += 1; pn(s, n)

s = vocab_write(prs, DAY, "现在", "Now",
                [("现", "xiàn", "8 笔", "「王」 + 「见」 — 现 在 看 见"),
                 ("在", "zài",  "6 笔", "上 「ナ」 下 「土」 — 站 在 那里")])
n += 1; pn(s, n)


# ============================================================
# 24 · PAST vs NOW MATCHING ACTIVITY
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔁 过去 vs 现在 · Past vs Now", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "找 一 找 — 古代 工具 ↔ 现 代 科技!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "Match: ancient tools ↔ modern tech",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

pairs = [
    ("🎋", "竹 简",      "📓", "笔记本",       "notebook"),
    ("📜", "手 写 书",   "📚", "印刷 书",      "printed books"),
    ("✉️", "写 信",      "📧", "email",        "email"),
    ("🗺️", "地 图",      "📍", "Google Maps",  "Google Maps"),
    ("🧮", "算 盘",      "🖩",  "计算 器",      "calculator"),
    ("🕯️", "蜡 烛",      "💡", "电 灯",        "electric light"),
]
pcol_w = 4.55; pcol_gap = 0.20
pcol_start = (10 - 2*pcol_w - pcol_gap)/2

for i, (em1, cn1, em2, cn2, en2) in enumerate(pairs):
    row = i % 3; col = i // 3
    x = pcol_start + col*(pcol_w + pcol_gap)
    y = 1.55 + row*1.08
    pc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(pcol_w), Inches(0.95))
    pc.fill.solid(); pc.fill.fore_color.rgb = WHITE
    pc.line.color.rgb = ANCIENT if col == 0 else DAY
    pc.line.width = Pt(2.5)
    tb(s, x+0.10, y+0.20, 0.55, 0.55, em1, sz=24, a=PP_ALIGN.CENTER)
    tb(s, x+0.70, y+0.30, 1.30, 0.45, cn1, sz=14, b=True, c=ANCIENT, a=PP_ALIGN.LEFT)
    tb(s, x+2.00, y+0.30, 0.40, 0.45, "↔", sz=22, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, x+2.45, y+0.20, 0.55, 0.55, em2, sz=24, a=PP_ALIGN.CENTER)
    tb(s, x+3.05, y+0.18, pcol_w-3.15, 0.35, cn2, sz=13, b=True, c=DAY, a=PP_ALIGN.LEFT)
    tb(s, x+3.05, y+0.52, pcol_w-3.15, 0.28, en2, sz=9, c=GRAY, a=PP_ALIGN.LEFT)

tb(s, 0.4, 5.10, 9.2, 0.28, "💡 古代 + 现代 — 解决 一 样 的 问题, 用 不 一 样 的 方法!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 配 对:\n• 让 学 生 一 个 一 个 想: 古 人 怎么 做? 现在 怎么 做?\n• 强 调: 「目 的 一 样 — 工具 不 同」")


# ============================================================
# 25 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3", "🎨 下午 3:00–4:30  ·  Project: 造纸 体 验!",
        DAY, "📜"); n += 1; pn(s, n)


# ============================================================
# 26 · COMPLETE TODAY'S BOOKLET (added per user request)
# ============================================================
s = booklet_slide(prs, day_num=4, day_topic_cn="科技 改变 生活 · 造纸 术",
                  day_color=DAY, page_count=4)
n += 1; pn(s, n)


# ============================================================
# 27 · PROJECT INTRO — 造纸 体验
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.4, 0.45), (9.1, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.40, 9.4, 0.40, "🏆 Project Time · 项 目 时 间!",
   sz=16, b=True, c=NEON, a=PP_ALIGN.CENTER)

tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.95), Inches(8.8), Inches(2.05))
tt.fill.solid(); tt.fill.fore_color.rgb = ANCIENT
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.20, 8.4, 0.85, "📜 造 纸 体 验!",
   sz=42, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.10, 8.4, 0.45, "Make Paper Like Cai Lun!",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.55, 8.4, 0.30, "今天 我们 自 己 做 一 张 纸 — 像 古 人 一 样!",
   sz=12, c=WARM, a=PP_ALIGN.CENTER)

si = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.20), Inches(8.8), Inches(0.70))
si.fill.solid(); si.fill.fore_color.rgb = WHITE
si.line.color.rgb = STAR; si.line.width = Pt(2)
tb(s, 0.75, 3.35, 8.5, 0.42, "👥 小 组 合作  ·  ⏱️ 60-90 分钟  ·  🎯 体 验 古代 中国 科技!",
   sz=13, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)

tb(s, 0.3, 4.20, 9.4, 0.95, "📰 → 💧 → 🍲 → 🟫 → 💪 → ☀️",
   sz=42, a=PP_ALIGN.CENTER)
tb(s, 0.3, 5.15, 9.4, 0.30, "撕 → 泡 → 浆 → 铺 → 压 → 晾",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "1-2 分钟 hype:\n• 大 声 宣 布: 「我 们 要 像 蔡 伦 一 样 做 纸!」\n• 让 学 生 兴 奋 起 来")


# ============================================================
# 28 · MATERIALS — 准备 工 具
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🛠️ 准备 工 具 · Materials", ANCIENT)

tb(s, 0.4, 0.85, 9.2, 0.30, "造 纸 需要 这 些 — 老师 提前 准 备 好!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

materials = [
    ("📰", "废 纸 / 旧 纸", "Waste paper"),
    ("💧", "水",          "Water"),
    ("🍲", "搅 拌 器", "Blender"),
    ("🪣", "盆 / 桶",      "Bowl / bucket"),
    ("🔲", "筛 网",        "Mesh screen"),
    ("🧽", "海 绵",        "Sponge"),
    ("📚", "毛 巾 / 旧 布", "Towel / cloth"),
    ("☀️", "晾 衣 架",     "Drying rack"),
]
mw = 2.20; mgap_x = 0.15; mgap_y = 0.18; mh = 1.45
mstart_x = (10 - 4*mw - 3*mgap_x)/2
for i, (em, cn, en) in enumerate(materials):
    row = i // 4; col = i % 4
    x = mstart_x + col*(mw + mgap_x)
    y = 1.40 + row*(mh + mgap_y)
    panel(s, x, y, mw, mh, ANCIENT, fill=WHITE, lw=2.5)
    tb(s, x, y+0.15, mw, 0.65, em, sz=34, a=PP_ALIGN.CENTER)
    tb(s, x, y+0.85, mw, 0.35, cn, sz=12, b=True, c=ANCIENT, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, y+1.18, mw-0.10, 0.25, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.55))
tip.fill.solid(); tip.fill.fore_color.rgb = DAY; tip.line.fill.background()
tb(s, 0.55, 4.95, 9.0, 0.32, "💡 老师 提示: 纸浆 提前 准 备 好 — 节 省 时 间!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.25, 9.0, 0.20, "Teacher tip: prepare pulp in advance",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟 介 绍 材料:\n• 主 要 材料 是 废 纸 — 回 收 利 用!\n• 围 裙 + 旧 衣 服 也 要 提 醒 学 生 穿\n• 准 备 几 桶 水 — 一 定 会 弄 湿 + 弄 脏\n• 报 纸 铺 桌 子 / 地 上 防 水")


# ============================================================
# 29 · 6-STEP PROCESS
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📋 6 步 造 纸 · Make Paper in 6 Steps", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "跟 着 这 6 步 — 你 也 能 像 蔡 伦 一 样!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

paper_steps = [
    ("1️⃣", "📰", "撕 废 纸",  "Tear waste paper", ANCIENT),
    ("2️⃣", "💧", "泡 水",     "Soak in water", CYBER),
    ("3️⃣", "🍲", "做 纸 浆",   "Make pulp", SEPIA),
    ("4️⃣", "🔲", "铺 在 筛 网", "Spread on screen", ORANGE),
    ("5️⃣", "💪", "压 水",      "Press out water", DAY),
    ("6️⃣", "☀️", "晾 干!",    "Dry it!", PURPLE),
]
sw = 2.95; sgap_x = 0.18; sgap_y = 0.20; sh = 1.55
sstart_x = (10 - 3*sw - 2*sgap_x)/2
for i, (num, em, cn, en, cl) in enumerate(paper_steps):
    row = i // 3; col = i % 3
    x = sstart_x + col*(sw + sgap_x)
    y = 1.40 + row*(sh + sgap_y)
    panel(s, x, y, sw, sh, cl, fill=WHITE, lw=2.5)
    tb(s, x+0.10, y+0.08, 0.55, 0.45, num, sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x+0.60, y+0.08, sw-0.70, 0.70, em, sz=38, a=PP_ALIGN.CENTER)
    tb(s, x, y+0.85, sw, 0.40, cn, sz=16, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, y+1.22, sw, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    if col < 2:
        arrow(s, x + sw + 0.02, y + sh/2 - 0.12, w=0.14, h=0.24, color=DAY)

tb(s, 0.4, 4.95, 9.2, 0.30, "✨ 不 完 美 也 没 关 系 — 重要 的 是 你 试 了!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.4, 5.25, 9.2, 0.22, "It's OK if not perfect — the experience matters!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 演示 + 30-45 分钟 操作:\n• 老师 先 演示 一 遍 整 个 流程\n• 然后 学 生 自 己 做\n• Step 3 (纸 浆) 老师 提前 准 备 — 节 省 时 间")


# ============================================================
# 30 · DISCUSSION — 反 思 + 引 导 问 题
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🤔 想 一 想 · Let's Reflect", DAY)

tb(s, 0.4, 0.85, 9.2, 0.32, "做 完 后 一 起 想 一 想!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

reflections = [
    ("❓", "为 什么 以前 的 人 需要 纸?",
     "Why did ancient people need paper?", ANCIENT),
    ("🎯", "纸 解决 了 什么 问题?",
     "What problem does paper solve?", DAY),
    ("🌍", "如果 没 有 纸 — 今天 的 生活 会 怎样?",
     "What if there were no paper today?", PURPLE),
]
rw = 4.40; rgap_x = 0.20; rh = 1.50
rstart_x = (10 - 2*rw - rgap_x)/2

for i in range(2):
    em, cn, en, cl = reflections[i]
    x = rstart_x + i*(rw + rgap_x)
    y = 1.30
    panel(s, x, y, rw, rh, cl, fill=WHITE, lw=3)
    tb(s, x+0.20, y+0.30, 0.85, 0.85, em, sz=40, a=PP_ALIGN.LEFT)
    tb(s, x+1.20, y+0.30, rw-1.35, 0.50, cn, sz=15, b=True, c=cl)
    tb(s, x+1.20, y+0.85, rw-1.35, 0.40, en, sz=10, c=GRAY)

em, cn, en, cl = reflections[2]
y = 1.30 + rh + 0.20
panel(s, rstart_x, y, 2*rw + rgap_x, rh, cl, fill=WHITE, lw=3)
tb(s, rstart_x+0.20, y+0.30, 0.85, 0.85, em, sz=40, a=PP_ALIGN.LEFT)
tb(s, rstart_x+1.20, y+0.35, 2*rw+rgap_x-1.35, 0.45, cn, sz=18, b=True, c=cl)
tb(s, rstart_x+1.20, y+0.85, 2*rw+rgap_x-1.35, 0.35, en, sz=11, c=GRAY)

tb(s, 0.4, 4.85, 9.2, 0.30, "💬 小 组 讨论 — 然后 一 起 分 享!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.4, 5.15, 9.2, 0.22, "Discuss in groups — then share with the class!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5-10 分钟 讨论:\n• 让 学 生 用 自 己 做 的 纸 当 道具\n• 引 导: 纸 帮 古 人 记 录 知识 / 没 纸 → 没 书\n• 联 系 上午: 「科技 = 解决 问题」")


# ============================================================
# 31 · SHOWCASE — 分 享 你 的 纸
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎤 分 享 你 的 纸 · Show Your Paper!", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "拿 起 你 做 的 纸 — 一 起 分 享!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

fp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.30), Inches(9.20), Inches(3.10))
fp.fill.solid(); fp.fill.fore_color.rgb = INK
fp.line.color.rgb = STAR; fp.line.width = Pt(3)
tb(s, 0.55, 1.42, 9.0, 0.35, "💬 用 句 型 分 享 — Share with these frames:",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)

share_frames = [
    ("1️⃣", "🛠️", "我 做 了 ______",
     "What I made: ___"),
    ("2️⃣", "💡", "纸 为 什么 重要? ______",
     "Why is paper important? ___"),
    ("3️⃣", "✨", "我 今天 学 到 了 ______",
     "Today I learned ___"),
]
for i, (num, em, cn, en) in enumerate(share_frames):
    y = 1.90 + i*0.78
    tb(s, 0.70, y, 0.60, 0.55, num, sz=22, b=True, c=STAR, a=PP_ALIGN.LEFT)
    tb(s, 1.35, y+0.05, 0.50, 0.45, em, sz=20, a=PP_ALIGN.LEFT)
    tb(s, 1.95, y+0.05, 7.55, 0.40, cn, sz=15, b=True, c=STAR, a=PP_ALIGN.LEFT)
    tb(s, 1.95, y+0.45, 7.55, 0.28, en, sz=9, c=WARM, a=PP_ALIGN.LEFT)

cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.60), Inches(9.20), Inches(0.85))
cb.fill.solid(); cb.fill.fore_color.rgb = ANCIENT
cb.line.color.rgb = STAR; cb.line.width = Pt(3)
tb(s, 0.55, 4.72, 9.0, 0.40, "🌟 今天 你 是 「小 蔡 伦」!",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.15, 9.0, 0.28, "Today you are a Little Cai Lun!",
   sz=11, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 分 享:\n• 每 个 学 生 30-60 秒\n• 拿 着 自 己 的 纸 介 绍\n• 鼓 励 中文 — 简 单 句 子 就 好\n• 拍 照 留 念 — 发 给 家 长!")


# ============================================================
# 32 · SHARE + CLOSE
# ============================================================
s = share_close(prs, DAY,
    frames_cn=["「以前 的 人 用 ______ 写 字」",
               "「纸 帮 我们 ______」"],
    frames_en="Ancient people wrote on ___  ·  Paper helps us ___",
    next_day_cn="Day 5 · 未来 科技 — 你 会 发明 什么?",
    next_day_en="Day 5 · Future tech — what will YOU invent?",
    next_emoji="🚀")
n += 1; pn(s, n)
notes(s, "5 分钟 收 尾:\n• 句 型 帮 学 生 巩 固 今天 学 的 — 纸 + 古代 科技\n• 明天 Day 5 — 未来 + 发明")


out = os.path.join(os.path.dirname(__file__), "day4_life.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
