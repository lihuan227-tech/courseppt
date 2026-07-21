#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
玩转创新科技 · Day 1: 认识 AI · 做聪明的 AI 小主人 (v4)
Based on the 《做聪明的AI小主人》 reference deck style + structured 3-session lesson plan.

Audience: K-5 mixed-age Chinese immersion summer camp (~20+ students)
Pedagogy: project-based, inquiry-based, voting games, scenario discussions
Design: heavy visuals, photo placeholders (NO auto-embedded resources), minimal text,
        teacher prompts on every slide via speaker notes
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = make_presentation()
DAY = AI_PURPLE
n = 0


# ============================================================
# Vote slide helper (used for AI-or-Not game + Smart-User scenarios)
# ============================================================
def vote_slide(s, title_cn, title_en, prompt_cn, item_visual_hint_cn, item_visual_hint_en,
               option_a_emoji, option_a_label, option_a_color,
               option_b_emoji, option_b_label, option_b_color,
               answer_correct_letter, answer_explanation_cn, answer_explanation_en):
    """Voting slide template: big visual placeholder + 2 voting options + reveal answer band."""
    bg(s, CREAM, prs)
    hb(s, title_cn + "  " + title_en, DAY)
    # Top prompt
    tb(s, 0.40, 0.85, 9.20, 0.30, prompt_cn, sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    # LEFT: large photo slot for the item
    photo_slot(s, 0.40, 1.25, 4.20, 3.20, item_visual_hint_cn, item_visual_hint_en, color=DAY)
    # RIGHT: two large vote buttons
    for i, (em, label, cl, letter) in enumerate([
        (option_a_emoji, option_a_label, option_a_color, "A"),
        (option_b_emoji, option_b_label, option_b_color, "B"),
    ]):
        y = 1.25 + i * 1.65
        is_correct = (letter == answer_correct_letter)
        panel_(s, 4.80, y, 4.80, 1.50, cl, lw=4 if is_correct else 2.5)
        tb(s, 4.95, y + 0.15, 0.85, 0.85, em, sz=44)
        tb(s, 5.85, y + 0.20, 3.60, 0.50, label, sz=20, b=True, c=cl)
        if is_correct:
            tb(s, 5.85, y + 0.75, 3.60, 0.55, "✨ 正确答案!", sz=14, b=True, c=cl)
    # Bottom: answer explanation band
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.90))
    band.fill.solid()
    band.fill.fore_color.rgb = INK
    band.line.color.rgb = STAR
    band.line.width = Pt(2)
    tb(s, 0.55, 4.72, 9.00, 0.30, "💡 答案解释  Answer", sz=10, b=True, c=STAR)
    tb(s, 0.55, 5.02, 9.00, 0.30, answer_explanation_cn, sz=12, b=True, c=WHITE)
    tb(s, 0.55, 5.32, 9.00, 0.22, answer_explanation_en, sz=8, c=WARM)
    return s


def panel_(s, l, t, w, h, color, fill=WHITE, lw=2.5):
    """Local thin wrapper around shared panel (kept for readability)."""
    return panel(s, l, t, w, h, color, fill=fill, lw=lw)


def scenario_slide(s, scene_num, scene_cn, scene_en, verdict_correct, reason_cn, reason_en):
    """Smart-AI-user scenario: one scene + 2 vote options + answer.
    verdict_correct = '✅' (smart) or '❌' (not smart)
    """
    bg(s, CREAM, prs)
    hb(s, f"🎯 情景 {scene_num} · 这样做, 对吗?  Scenario {scene_num}", DAY)
    # Scene card — big quote bubble
    bubble = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.40), Inches(1.00), Inches(9.20), Inches(1.65))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = WARM
    bubble.line.color.rgb = DAY
    bubble.line.width = Pt(2.5)
    tb(s, 0.55, 1.20, 9.00, 0.55, f"「{scene_cn}」", sz=20, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 1.85, 9.00, 0.40, scene_en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 2.25, 9.00, 0.30, "👉 这算聪明的 AI 小主人吗?",
       sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    # 2 vote buttons
    for i, (em, label, cl, letter) in enumerate([
        ("✅", "聪明小主人  Smart!", GREEN, "✅"),
        ("❌", "不聪明  Not smart", PINK, "❌"),
    ]):
        x = 0.40 + i * 4.80
        is_correct = (letter == verdict_correct)
        panel_(s, x, 2.85, 4.40, 1.10, cl, lw=4 if is_correct else 2.5)
        tb(s, x + 0.15, 3.00, 0.85, 0.80, em, sz=36)
        tb(s, x + 1.00, 3.05, 3.30, 0.45, label, sz=15, b=True, c=cl)
        if is_correct:
            tb(s, x + 1.00, 3.50, 3.30, 0.40, "✨ 答案!", sz=13, b=True, c=cl)
    # Answer reasoning
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.40), Inches(4.10), Inches(9.20), Inches(1.45))
    band.fill.solid()
    band.fill.fore_color.rgb = INK
    band.line.color.rgb = STAR
    band.line.width = Pt(2)
    tb(s, 0.55, 4.20, 9.00, 0.30, "💡 为什么?  Why?", sz=11, b=True, c=STAR)
    tb(s, 0.55, 4.55, 9.00, 0.65, reason_cn, sz=12, b=True, c=WHITE)
    tb(s, 0.55, 5.20, 9.00, 0.30, reason_en, sz=9, c=WARM)
    return s


# ============================================================
# 1 · COVER
# ============================================================
s = ns(prs); bg(s, INK, prs)

tb(s, 0.5, 0.35, 9.0, 0.40, "🚀 玩转创新科技  Playing with Innovative Tech",
   sz=16, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.5, 0.85, 9.0, 0.65, "Day 1 · 认识 AI", sz=38, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.5, 1.55, 9.0, 0.50, "做聪明的 AI 小主人", sz=24, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.5, 2.10, 9.0, 0.35, "Becoming a Smart AI Master · Getting to Know AI",
   sz=13, c=LGRAY, a=PP_ALIGN.CENTER)
# Mascot placeholders (rounded oval shapes)
photo_slot(s, 2.30, 2.75, 2.50, 2.20,
           "AI 小主人形象 / 卡通机器人",
           "Cute AI mascot or robot character", color=DAY)
photo_slot(s, 5.20, 2.75, 2.50, 2.20,
           "小朋友用 AI 的照片",
           "Kids using AI", color=ORANGE)
tb(s, 0.5, 5.10, 9.0, 0.35, "✨ 一起出发 — AI 探索之旅!", sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Day 1 节奏:\n• Session 1 (45 min): Warm-up + AI/not-AI游戏 + AI犯错 + 聪明小主人\n• Session 2 (45 min): 复习 + 词汇 + Mini Booklet\n• Session 3 (90 min): AI 故事书项目 + 选做 AI 机器人手工\n\n班级: 20+ K-5 immersion 学生\n要点: 全中文 + 大量图片 + 多互动")

# ============================================================
# 2 · AI 探索地图 (学习目标)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🗺️ AI 探索地图 · 今天的学习目标  Our Learning Map", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "今天我们要一起完成 5 件事 — 你准备好了吗?",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

goals = [
    ("1️⃣", "🤖", "知道 AI 是帮人解决问题的聪明工具", "Understand AI = smart tool that helps", DAY),
    ("2️⃣", "🌍", "认识生活中常见的 AI 应用", "Recognize AI in daily life", CYBER),
    ("3️⃣", "🔍", "知道 AI 也会犯错 — 不是什么都对", "Know AI makes mistakes too", PINK),
    ("4️⃣", "🛡️", "学会安全 + 负责地用 AI", "Use AI safely & responsibly", GREEN),
    ("5️⃣", "🎨", "设计自己的 AI 小帮手", "Design your own AI helper", ORANGE),
]
for i, (num, em, cn, en, cl) in enumerate(goals):
    y = 1.30 + i * 0.78
    panel(s, 0.50, y, 9.00, 0.65, cl, fill=WHITE, lw=2)
    tb(s, 0.60, y + 0.13, 0.55, 0.45, num, sz=18, b=True, c=cl)
    tb(s, 1.15, y + 0.13, 0.55, 0.45, em, sz=20)
    tb(s, 1.75, y + 0.10, 5.80, 0.35, cn, sz=12, b=True, c=DARK)
    tb(s, 1.75, y + 0.40, 5.80, 0.25, en, sz=9, c=GRAY)
n += 1; pn(s, n)

# ============================================================
# 3 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1", "🌅 上午 11:00–11:45  ·  认识 AI + 做聪明的 AI 小主人", DAY, "🤖")
n += 1; pn(s, n)

# ============================================================
# 4 · WARM-UP — 你用过 AI 吗?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "👋 热身 · 你用过 AI 吗?  Have You Used AI?", DAY)
tb(s, 0.40, 0.85, 9.20, 0.28, "举手! Raise your hand if you've ___ ",
   sz=13, b=True, c=DAY, a=PP_ALIGN.CENTER)

# 5 questions as numbered rows
questions = [
    ("🎙️", "谁和 Siri / Alexa 说过话?", "Talked to Siri / Alexa?"),
    ("🤖", "谁家有扫地机器人?", "Has a Roomba at home?"),
    ("🗺️", "谁用过地图导航 (Google Maps)?", "Used GPS navigation?"),
    ("🌐", "谁用过翻译软件?", "Used a translation app?"),
    ("🤖", "谁见过真的机器人?", "Has seen a real robot?"),
]
for i, (em, cn, en) in enumerate(questions):
    y = 1.25 + i * 0.72
    panel(s, 0.50, y, 9.00, 0.62, DAY, fill=WHITE, lw=2)
    tb(s, 0.65, y + 0.10, 0.50, 0.45, em, sz=22)
    tb(s, 1.25, y + 0.08, 6.50, 0.32, cn, sz=14, b=True, c=DARK)
    tb(s, 1.25, y + 0.36, 6.50, 0.26, en, sz=9, c=GRAY)
    # Hand-raise indicator on right
    tb(s, 7.95, y + 0.12, 1.45, 0.45, "✋ 举手!", sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)

# Activity callout
activity_box(s, 0.40, 4.85, 9.20, 0.55,
             "数一数 — 哪一个举手的人最多?",
             "Count — which one had the most hands?",
             color=DAY)
n += 1; pn(s, n)
notes(s, "5 分钟:\n• 老师一个一个念问题, 看哪个举手多\n• 鼓励学生喊出自己用过的 AI\n• 不用全答对 — 重点是引出兴趣")

# ============================================================
# 5 · Turn & Talk + AI 定义 summary
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "💬 Turn & Talk · 你觉得 AI 是什么?", DAY)

# Top: Turn & Talk prompt
panel(s, 0.40, 0.95, 9.20, 1.30, DAY, fill=WARM)
tb(s, 0.55, 1.05, 9.00, 0.40, "👯 转身! 告诉同桌:",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.50, 9.00, 0.50, "「你觉得 AI 是什么?」",
   sz=22, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.02, 9.00, 0.28, "Turn to your partner: 'What do YOU think AI is?'  (1 分钟)",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom: Teacher summary card
panel(s, 0.40, 2.50, 9.20, 2.50, CYBER, fill=WHITE)
panel_head(s, 0.40, 2.50, 9.20, CYBER, "✨ 老师总结  Teacher's Summary", sz=12)
tb(s, 0.55, 3.10, 9.00, 0.45, "「AI 就像一个聪明的小助手!」",
   sz=22, b=True, c=CYBER, a=PP_ALIGN.CENTER)
tb(s, 0.55, 3.60, 9.00, 0.30, "AI is like a smart little helper!",
   sz=11, c=GRAY, a=PP_ALIGN.CENTER)
# 4 abilities as small chips
abilities = [
    ("👂", "能听", "Listens"),
    ("👀", "能看", "Sees"),
    ("💬", "能答问题", "Answers Qs"),
    ("🛠️", "帮我们解决问题", "Solves problems"),
]
chip_w = 2.10; gap = 0.08
total = 4 * chip_w + 3 * gap; start = (10 - total) / 2
for i, (em, cn, en) in enumerate(abilities):
    x = start + i * (chip_w + gap)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(4.05), Inches(chip_w), Inches(0.90))
    chip.fill.solid(); chip.fill.fore_color.rgb = WARM
    chip.line.color.rgb = CYBER; chip.line.width = Pt(1.5)
    tb(s, x, 4.10, chip_w, 0.35, em, sz=18, a=PP_ALIGN.CENTER)
    tb(s, x, 4.45, chip_w, 0.30, cn, sz=11, b=True, c=CYBER, a=PP_ALIGN.CENTER)
    tb(s, x, 4.72, chip_w, 0.22, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟:\n• 1 分钟同桌互相说\n• 2-3 个学生全班分享\n• 老师念 summary, 全班跟读\n• 鼓励学生用自己的话解释")

# ============================================================
# 6 · AI or Not AI? Game intro + criteria
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎯 游戏时间 · 是 AI 吗?  Is It AI?", ORANGE)

# Intro
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.80))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = ORANGE; intro.line.width = Pt(2)
tb(s, 0.55, 1.05, 9.00, 0.30, "⚠️ 注意! 不是所有电子产品都是 AI",
   sz=14, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.40, 9.00, 0.28, "Not everything electronic is AI!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Two-column criteria
panel(s, 0.40, 1.95, 4.55, 3.10, GREEN)
panel_head(s, 0.40, 1.95, 4.55, GREEN, "✅ 这些算 AI", sz=13)
ai_can = [
    "👂 听懂 / 看懂信息",
    "🧠 自己做判断",
    "💬 回答问题",
    "🎯 给建议",
    "🛠️ 帮你解决问题",
]
for i, line in enumerate(ai_can):
    tb(s, 0.55, 2.55 + i * 0.45, 4.30, 0.40, line, sz=12, b=True, c=DARK)

panel(s, 5.05, 1.95, 4.55, 3.10, PINK)
panel_head(s, 5.05, 1.95, 4.55, PINK, "❌ 这些不算 AI", sz=13)
ai_not = [
    "💡 只会 「开 / 关」",
    "🧮 固定计算 (1+1=2)",
    "🎵 固定动作 (一直转)",
    "📺 只播你选的节目",
    "⚙️ 没有 「思考」",
]
for i, line in enumerate(ai_not):
    tb(s, 5.20, 2.55 + i * 0.45, 4.30, 0.40, line, sz=12, b=True, c=DARK)

activity_box(s, 0.40, 5.10, 9.20, 0.50,
             "🎮 准备好 — 接下来 7 个物品, 一起投票!",
             "Get ready — 7 items, you vote!", color=ORANGE)
n += 1; pn(s, n)
notes(s, "3-5 分钟:\n• 强调 「不是所有电子产品都是 AI」\n• 用简单标准: 能 「思考」 + 「学习」 才是 AI\n• 跟读几个关键词")

# ============================================================
# 7-11 · AI or Not AI? — 5 Contrast Pairs (普通 vs AI version)
# ============================================================
def contrast_slide(s, title_cn, title_en,
                   left_label_cn, left_label_en, left_photo_cn, left_photo_en, left_desc_cn,
                   right_label_cn, right_label_en, right_photo_cn, right_photo_en, right_desc_cn,
                   key_diff_cn, key_diff_en):
    """Contrast voting slide: 普通 X vs AI X side by side + key difference reveal."""
    bg(s, CREAM, prs)
    hb(s, title_cn + "  " + title_en, DAY)
    tb(s, 0.40, 0.85, 9.20, 0.30, "🤔 哪个 是 AI? 它们 有 什么 不一样?",
       sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 1.13, 9.20, 0.22, "Which one is AI? What's the difference?",
       sz=9, c=GRAY, a=PP_ALIGN.CENTER)

    # LEFT card — 普通 version (NOT AI)
    panel(s, 0.40, 1.45, 4.55, 2.90, PINK)
    # Header strip
    head_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.40), Inches(1.45), Inches(4.55), Inches(0.42))
    head_l.fill.solid(); head_l.fill.fore_color.rgb = PINK; head_l.line.fill.background()
    tb(s, 0.40, 1.50, 4.55, 0.32, f"❌ {left_label_cn}  {left_label_en}",
       sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Photo placeholder
    photo_slot(s, 0.55, 1.97, 4.25, 1.65, left_photo_cn, left_photo_en, color=PINK)
    # Description below photo
    tb(s, 0.55, 3.70, 4.25, 0.60, left_desc_cn, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

    # RIGHT card — 智能 version (IS AI)
    panel(s, 5.05, 1.45, 4.55, 2.90, GREEN)
    head_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(5.05), Inches(1.45), Inches(4.55), Inches(0.42))
    head_r.fill.solid(); head_r.fill.fore_color.rgb = GREEN; head_r.line.fill.background()
    tb(s, 5.05, 1.50, 4.55, 0.32, f"✅ {right_label_cn}  {right_label_en}",
       sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    photo_slot(s, 5.20, 1.97, 4.25, 1.65, right_photo_cn, right_photo_en, color=GREEN)
    tb(s, 5.20, 3.70, 4.25, 0.60, right_desc_cn, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

    # Key difference reveal at bottom
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.40), Inches(4.50), Inches(9.20), Inches(1.05))
    band.fill.solid(); band.fill.fore_color.rgb = INK
    band.line.color.rgb = STAR; band.line.width = Pt(2)
    tb(s, 0.55, 4.58, 9.00, 0.28, "💡 区别 在 哪里?  Key Difference:",
       sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 4.88, 9.00, 0.36, key_diff_cn, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 5.25, 9.00, 0.22, key_diff_en, sz=9, c=WARM, a=PP_ALIGN.CENTER)
    return s


contrasts = [
    # 1. Lights
    ("💡 灯 · 对比",  "Lights · Compare",
     "普通 电灯", "Regular Light Bulb",
     "普通 灯泡 + 墙上 开关", "Light bulb + wall switch on/off",
     "只 会 「开 / 关」",
     "智能 灯", "Smart Light",
     "Alexa/Hue 智能 灯 + 手机 App", "Smart bulb controlled by voice/app",
     "「Alexa, 开 灯!」 — 它 听懂!",
     "AI 能 听 懂 你 的 话!  —  普通 灯 只 等 你 按 开关",
     "AI understands your voice — a regular light just waits for the switch"),

    # 2. Calculator vs Siri
    ("🧮 算术 · 对比",  "Math · Compare",
     "普通 计算器", "Regular Calculator",
     "普通 计算器 (按 按钮)", "Plastic calculator with buttons",
     "你 按 1 + 1, 它 答 2 — 不 会 听 话",
     "Siri / Alexa", "Voice Assistant",
     "Siri 在 手机 屏幕 上 听 你 问 题", "Siri on phone responding to voice",
     "你 说 「1 加 1 等于 几?」 — 它 答!",
     "AI 能 听 懂 你 说 的 话  —  计算器 只 认 按钮",
     "AI understands speech — calculators only know buttons"),

    # 3. Cars
    ("🚗 车 · 对比",  "Cars · Compare",
     "普通 汽车", "Regular Car",
     "普通 汽车 + 司机 开 车", "Regular car with human driver",
     "人 看 路, 人 开 车",
     "自动 驾驶 车", "Self-Driving Car",
     "Tesla / Waymo 自动 驾驶 车", "Tesla or Waymo car driving itself",
     "车 自己 看 路, 自己 判断!",
     "AI 能 自己 「看」 + 「判 断」  —  普通 车 要 人 来 开",
     "AI can SEE and DECIDE on its own — regular cars need a human"),

    # 4. Toys
    ("🧸 玩具 · 对比",  "Toys · Compare",
     "普通 玩具 熊", "Plush Bear",
     "毛绒 玩具 熊", "Stuffed teddy bear",
     "你 问 什么, 它 都 不 答",
     "智能 玩具", "AI Talking Toy",
     "会 说话 的 AI 玩具 (如 Moxie)", "AI chatbot toy that talks",
     "你 问, 它 答 + 还 会 学!",
     "AI 能 聊 天 + 还 会 学 习  —  普通 玩具 只 是 软软 的",
     "AI toys chat + learn — regular toys just sit there"),

    # 5. Maps
    ("🗺️ 地图 · 对比",  "Maps · Compare",
     "纸 地图", "Paper Map",
     "一 张 纸 地图 / 旧 式 地图", "Old printed paper map",
     "你 自己 看, 自己 找 路",
     "Google Maps", "GPS / Google Maps",
     "Google Maps 手机 App 截图", "Google Maps on phone with route",
     "推荐 路线 + 避 开 堵 车!",
     "AI 能 「想」 哪 条 路 最 快  —  纸 地图 不 会 变",
     "AI thinks about the best route — paper maps stay still"),
]
for c in contrasts:
    s = ns(prs)
    contrast_slide(s, *c)
    n += 1; pn(s, n)

# ============================================================
# 14 · AI 也会犯错! (3 REAL documented AI hallucinations)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🤔 AI 也会犯错!  AI Can Be Wrong Too!", PINK)
tb(s, 0.40, 0.85, 9.20, 0.28, "这 3 个 都 是 真实 发生过 的 AI 错误 — 你能 找出 哪里 不对 吗?",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.40, 1.13, 9.20, 0.22, "These are REAL AI mistakes — can you spot what's wrong?",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# 3 REAL documented AI mistakes (kid-verifiable)
# Each: emoji, AI claim CN, AI claim EN, truth/why-wrong, accent color
wrongs = [
    ("🍓",
     "「Strawberry 这 个 词 里 有 2 个 R」",
     "'The word strawberry has 2 R's.'",
     "✋ 自己 数! S-t-R-a-w-b-e-R-R-y → 3 个 R!",
     "Count yourself — 3 R's!",
     CYBER),
    ("🍕",
     "「Pizza 涂 一点 胶水, 奶酪 不会 掉。」",
     "'Put a little glue on pizza so cheese stays.'",
     "⚠️ 真实事件! 2024 年 Google AI 出过 这 错 — 胶水 不能 吃!",
     "Real! Google AI suggested this in 2024 — never eat glue!",
     ORANGE),
    ("🐋",
     "「鲸鱼 是 一 种 鱼。」",
     "'Whales are a kind of fish.'",
     "🧠 实际: 鲸鱼 是 「哺乳 动物」 — 像 我们 一样!",
     "Actually: whales are mammals — like us!",
     GREEN),
]
card_w = 2.95; gap = 0.12
total = 3 * card_w + 2 * gap; start = (10 - total) / 2
for i, (em, cn, en, truth_cn, truth_en, cl) in enumerate(wrongs):
    x = start + i * (card_w + gap)
    panel(s, x, 1.45, card_w, 3.25, PINK, lw=2.5)
    # Top: "🤖 AI 说" label
    head_label = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x + 0.10), Inches(1.55), Inches(card_w - 0.20), Inches(0.30))
    head_label.fill.solid(); head_label.fill.fore_color.rgb = PINK; head_label.line.fill.background()
    tb(s, x + 0.10, 1.58, card_w - 0.20, 0.25, "🤖 AI 说:", sz=10, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Emoji
    tb(s, x, 1.92, card_w, 0.55, em, sz=32, a=PP_ALIGN.CENTER)
    # AI claim
    tb(s, x + 0.12, 2.50, card_w - 0.24, 0.65, cn, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x + 0.12, 3.15, card_w - 0.24, 0.30, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Divider line
    tb(s, x, 3.48, card_w, 0.20, "─ ─ ─", sz=10, c=PINK, a=PP_ALIGN.CENTER)
    # Truth reveal
    tb(s, x + 0.12, 3.68, card_w - 0.24, 0.65, truth_cn, sz=9, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.12, 4.30, card_w - 0.24, 0.32, truth_en, sz=7, c=GRAY, a=PP_ALIGN.CENTER)

# Concept reveal
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.75))
band.fill.solid(); band.fill.fore_color.rgb = INK
band.line.color.rgb = STAR; band.line.width = Pt(2)
tb(s, 0.55, 4.92, 9.00, 0.30, "💡 这 叫 「AI 幻觉」 · AI 会 「一本正经 地 胡说 八道」!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.25, 9.00, 0.30, "This is called 'AI Hallucination' — we must learn to check AI's answers!",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟讨论:\n\n这 3 个 都 是 真实 AI 出过 的 错误:\n\n1️⃣ STRAWBERRY 字母:\n• ChatGPT, Claude 等 早期 经常 答 错\n• 因为 AI 不 是 「真 看」 字母, 是 「学」 出 来 的\n• 让 学生 一起 数 — 全班 喊 R 的 位置\n• 教训: AI 不擅长 「精确 数 数」\n\n2️⃣ PIZZA 胶水 (真实 + 重要!):\n• 2024 年 5 月, Google 新出 「AI Overview」 推荐 把 「无毒 胶水」 放在 pizza 上 防止 奶酪 掉\n• 全球 上 了 新闻!\n• 原因: AI 学了 Reddit 上 一个 玩笑 帖子\n• 教训: AI 不 知道 什么 安全, 不要 听 AI 关于 食物 / 安全 的 建议\n\n3️⃣ 鲸鱼 = 鱼?\n• AI 看 「外表 像」 容易 混淆\n• 鲸鱼: 哺乳 动物 (mammal) — 用 肺 呼吸, 喂奶, 温血\n• 鱼: 用 鳃 呼吸, 冷血\n• K-5 应该 学过 mammal 概念\n• 教训: AI 凭 「样子」 容易 答 错\n\n讨论 引导:\n• 「为什么 AI 会 答 错?」 → 因为 它 学过 错的 信息 / 不擅长 精确\n• 「我们 怎么 知道 AI 错 了?」 → 自己 验证, 多问 几个 来源\n• 引出 下 一 张 slide: 小侦探 三步法")

# ============================================================
# 15 · AI 猜动物 — Talk-based AI guessing (moved up before Quick Draw)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🐾 AI 猜猜看 · 你描述, AI 来猜!", GREEN)

# Top intro
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.65))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = GREEN; intro.line.width = Pt(2)
tb(s, 0.55, 1.02, 9.00, 0.30, "🎯 心里想一个动物 → 给 AI 3 个提示 → AI 猜!",
   sz=13, b=True, c=GREEN, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.32, 9.00, 0.24, "Think of an animal → give 3 hints → AI guesses!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: 3-step process
panel(s, 0.40, 1.75, 4.55, 3.20, GREEN)
panel_head(s, 0.40, 1.75, 4.55, GREEN, "📝 怎么玩  How to Play", sz=12)
steps = [
    ("1️⃣", "🤔", "心里想一个动物", "Think of an animal"),
    ("2️⃣", "💬", "给 AI 3 个提示", "Give AI 3 hints"),
    ("3️⃣", "🤖", "AI 猜! 对了 = 赢!", "AI guesses — right = win!"),
]
for i, (num, em, cn, en) in enumerate(steps):
    y = 2.30 + i * 0.75
    tb(s, 0.55, y, 0.45, 0.45, num, sz=18, b=True, c=GREEN)
    tb(s, 1.05, y, 0.50, 0.45, em, sz=22)
    tb(s, 1.65, y + 0.02, 3.20, 0.35, cn, sz=12, b=True, c=DARK)
    tb(s, 1.65, y + 0.40, 3.20, 0.28, en, sz=8, c=GRAY)

# RIGHT: example dialogue
panel(s, 5.05, 1.75, 4.55, 3.20, ORANGE)
panel_head(s, 5.05, 1.75, 4.55, ORANGE, "💡 示例 · 猜动物  Example", sz=12)
hints = [
    ("👤 我:", "它有胡子...", CYBER),
    ("👤 我:", "它会喵喵叫...", CYBER),
    ("👤 我:", "它喜欢吃鱼!", CYBER),
    ("🤖 AI:", "是猫吗?", AI_PURPLE),
]
for i, (who, line, color) in enumerate(hints):
    y = 2.30 + i * 0.55
    tb(s, 5.20, y, 0.75, 0.30, who, sz=11, b=True, c=color)
    tb(s, 6.00, y + 0.02, 3.50, 0.32, line, sz=13, b=True, c=DARK)

# Bottom reflection band
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.55))
band.fill.solid(); band.fill.fore_color.rgb = INK
band.line.color.rgb = STAR; band.line.width = Pt(2)
tb(s, 0.55, 5.12, 9.00, 0.28, "💡 想一想: AI 真的「知道」吗? 还是像「背了很多答案的百科全书」?",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.40, 9.00, 0.18, "Does AI really 'know'? Or is it a giant encyclopedia of memorized answers?",
   sz=7, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟玩法:\n\n准备:\n• 老师打开 ChatGPT / Claude / Siri (用投影)\n• 准备「动物词库」给学生参考 (猫/狗/大象/长颈鹿/企鹅...)\n\nRound 1 (老师示范) — 3 分钟:\n• 老师心里选一个动物\n• 给 AI 念 3 个提示, 例: 「它很大」「有长鼻子」「在非洲」\n• AI 猜 → 大象!\n• 让学生看 AI 怎么「想」\n\nRound 2 (学生轮流) — 5-7 分钟:\n• 老师抽 3-4 个学生上来\n• 每人选一个动物, 用中文给 3 个提示\n• 全班先猜, 然后老师把提示输入给 AI\n• 看 AI 和学生谁猜得对\n\n反思引导 (重要!):\n• 「AI 是真的知道你想什么吗?」 → 不!\n• 「那 AI 怎么猜对的?」 → 它学过很多动物的特征\n• 「像不像背了很多答案的百科全书?」 → 像!\n• 关键比喻: AI = 「会查百科的电脑朋友」\n\nK-2 简化:\n• 用图片提示 (老师拿动物图片让学生描述)\n• 只给 2 个提示\n\n3-5 加深:\n• 给「难」的动物 (海豚? 蝴蝶?)")

# ============================================================
# 16 · AI 猜猜看 升级 — Student writes prompt, AI guesses (NEW)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎯 AI 猜猜看 升级版 · 你写, AI 来猜!", PURPLE)

# Top intro
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.70))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = PURPLE; intro.line.width = Pt(2)
tb(s, 0.55, 1.02, 9.00, 0.32, "🚀 升级! 不 只是 动物 — 任何 东西 都 可以 让 AI 猜!",
   sz=13, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.38, 9.00, 0.24, "Upgrade! Not just animals — ANY object, place, food, or character!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: 4-step process
panel(s, 0.40, 1.78, 4.55, 3.20, PURPLE)
panel_head(s, 0.40, 1.78, 4.55, PURPLE, "📝 怎么 玩  How to Play", sz=12)
steps = [
    ("1️⃣", "💭", "心里 想 一 个 东西", "Think of anything"),
    ("2️⃣", "✏️", "写 4 个 中文 提示", "Write 4 Chinese hints"),
    ("3️⃣", "🤖", "AI 猜! 看 几 次 对", "AI guesses — how many tries?"),
    ("4️⃣", "🎯", "对 了 = 你 的 提示 写 得 好!", "Got it = your hints were clear!"),
]
for i, (num, em, cn, en) in enumerate(steps):
    y = 2.30 + i * 0.62
    tb(s, 0.55, y, 0.45, 0.40, num, sz=16, b=True, c=PURPLE)
    tb(s, 1.05, y, 0.50, 0.40, em, sz=20)
    tb(s, 1.65, y + 0.02, 3.20, 0.32, cn, sz=11, b=True, c=DARK)
    tb(s, 1.65, y + 0.35, 3.20, 0.26, en, sz=8, c=GRAY)

# RIGHT: 5 category cards (what kinds of things to think of)
panel(s, 5.05, 1.78, 4.55, 3.20, ORANGE)
panel_head(s, 5.05, 1.78, 4.55, ORANGE, "💡 任何 东西 都 可以  Anything!", sz=12)
categories = [
    ("🍎", "物品", "苹果 / 雨伞 / 笔"),
    ("🏰", "地点", "公园 / 海边 / 学校"),
    ("👨‍🚀", "人物", "医生 / 老师 / 警察"),
    ("🍕", "食物", "披萨 / 包子 / 饺子"),
    ("🦸", "卡通", "皮卡丘 / 米奇 / 海绵宝宝"),
]
for i, (em, cn_t, cn_eg) in enumerate(categories):
    y = 2.30 + i * 0.52
    tb(s, 5.20, y, 0.45, 0.40, em, sz=18)
    tb(s, 5.75, y + 0.04, 1.10, 0.32, cn_t, sz=12, b=True, c=ORANGE)
    tb(s, 6.95, y + 0.06, 2.60, 0.30, cn_eg, sz=9, c=DARK)

# Bottom example/tip
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.50))
band.fill.solid(); band.fill.fore_color.rgb = INK
band.line.color.rgb = STAR; band.line.width = Pt(2)
tb(s, 0.55, 5.16, 9.00, 0.28, "✨ 示例: 「红色 + 圆 + 树上长 + 甜」 → AI: 是 苹果 吗?",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.42, 9.00, 0.18, "Example: 'Red + round + grows on tree + sweet' → AI: an apple?",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10-12 分钟升级版玩法 (紧接上一活动):\n\n升级点 (和 上一 张 对比):\n• 上一张: 只 动物 + 3 个提示\n• 这一张: 任何东西 + 4 个提示 + 学生 自己 写\n\nRound 1 (老师示范) — 3 分钟:\n• 老师想一个东西, 例如「公园」\n• 老师写 4 个提示: 「有滑梯」「有秋千」「小朋友 在 玩」「在 户外」\n• 输入 AI → AI 猜: 公园!\n• 让学生看 AI 怎么 「拼凑」 出答案\n\nRound 2 (学生 — 关键活动!) — 7-9 分钟:\n• 每个学生 在 纸 上 写 4 个 中文 提示 (心里 想 一 个 东西)\n• 抽 4-5 个学生 来 全班 念 自己 的 提示\n• 老师 把 提示 念 给 AI\n• 全班 + AI 一起 猜!\n• 看 谁 写 的 提示 让 AI 第 1 次 就 猜 对\n\n反思 (重要!):\n• 「为什么 有的 学生 的 提示, AI 一 下 就 猜 到? 有的 半 天 猜 不 到?」\n• 引导: 提示 越 清楚 + 越 具体 → AI 越 容易 猜\n• 这是 「写 prompt 的 技巧」 — 给 Day 3 ML 和 Session 3 项目 铺垫\n\n分层:\n• K-2: 老师帮写提示, 学生口头说\n• 3-5: 自己写, 写 4-5 个提示, 用 「形状/颜色/用途/在哪儿」 4 维度\n\n衔接 Quick Draw (下一活动):\n• 这一活动: 用 「文字」 告诉 AI\n• 下一活动 (Quick Draw): 用 「图画」 告诉 AI\n• AI 都能 「猜」 — 但 它真的 「懂」 吗?")

# ============================================================
# 17 · Quick Draw 挑战 — Live AI experience
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎨 Quick Draw · AI 猜猜看!  Can AI Guess?", ORANGE)

# Top: URL + setup
url_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.65))
url_bar.fill.solid(); url_bar.fill.fore_color.rgb = ORANGE; url_bar.line.fill.background()
tb(s, 0.55, 1.02, 9.00, 0.30, "🌐 网址: quickdraw.withgoogle.com  ·  免费, 不用注册!",
   sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.32, 9.00, 0.24, "AI 给任务 (画苹果!) — 你画 — AI 猜!",
   sz=10, b=True, c=STAR, a=PP_ALIGN.CENTER)

# 3 rounds
rounds = [
    ("1️⃣", "👩‍🏫", "Round 1", "老师 Demo",
     "老师画 「苹果」", "AI 猜几次才对?",
     "Teacher draws an apple — how many guesses?", CYBER),
    ("2️⃣", "🙋", "Round 2", "学生挑战",
     "选 2-3 个学生来画 (20 秒!)", "全班先猜: AI 会对吗?",
     "2-3 student volunteers (20 sec each)", PURPLE),
    ("3️⃣", "🎭", "Round 3", "难倒 AI!",
     "画难的: 🚲 / 🐉 / ✈️", "故意画 「不像」 — AI 还能猜吗?",
     "Try to fool AI: bike, dragon, plane", PINK),
]
card_w = 2.95; gap = 0.12
total = 3 * card_w + 2 * gap; start = (10 - total) / 2
for i, (num, em, rd, cn_t, cn_what, cn_q, en_d, cl) in enumerate(rounds):
    x = start + i * (card_w + gap)
    panel(s, x, 1.85, card_w, 2.80, cl, lw=2.5)
    # Header
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(1.85), Inches(card_w), Inches(0.45))
    head.fill.solid(); head.fill.fore_color.rgb = cl; head.line.fill.background()
    tb(s, x + 0.10, 1.91, 0.50, 0.32, num, sz=14, b=True, c=WHITE)
    tb(s, x + 0.60, 1.91, card_w - 0.70, 0.32, rd, sz=13, b=True, c=WHITE)
    # Icon
    tb(s, x, 2.40, card_w, 0.70, em, sz=40, a=PP_ALIGN.CENTER)
    # Title
    tb(s, x + 0.10, 3.10, card_w - 0.20, 0.35, cn_t, sz=14, b=True, c=cl, a=PP_ALIGN.CENTER)
    # What to do
    tb(s, x + 0.15, 3.50, card_w - 0.30, 0.42, cn_what, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)
    # Question prompt
    tb(s, x + 0.15, 3.92, card_w - 0.30, 0.42, cn_q, sz=10, b=True, c=cl, a=PP_ALIGN.CENTER)
    # English
    tb(s, x + 0.15, 4.35, card_w - 0.30, 0.28, en_d, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Activity flow at bottom
activity_box(s, 0.40, 4.80, 9.20, 0.75,
             "🗳️ 每轮之前: 全班投票 — 「AI 会猜对吗?」 👍 / 👎",
             "Before each round: class votes — Will AI guess right? 👍 / 👎",
             color=ORANGE)
n += 1; pn(s, n)
notes(s, "15-20 分钟互动:\n• 网址: quickdraw.withgoogle.com\n• 老师用投影 + iPad / 笔记本 (鼠标或触摸都可以)\n\nRound 1 (Demo) — 5 分钟:\n• 老师念 「Draw an apple」, 故意画慢一点\n• 全班看 AI 实时猜: 「Is it a ball? Is it a tomato? Is it an apple!」\n• 暂停问学生: 「AI 为什么先猜错?」\n\nRound 2 (Student) — 5-8 分钟:\n• 选 2-3 个学生上来 (举手最高的)\n• 20 秒限时\n• 全班先投票\n\nRound 3 (Trick AI) — 5 分钟:\n• 故意画难的 / 抽象的\n• 看 AI 哪里出错\n• 引出 「AI 是在比较以前看过的图」\n\n关键: 让学生体验, 不是看老师玩!")

# ============================================================
# 18 · Quick Draw 反思 — Why does AI guess wrong?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🤔 Quick Draw 反思 · AI 真的 「懂」 吗?", PURPLE)

# Top discussion prompt
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.70))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = PURPLE; intro.line.width = Pt(2)
tb(s, 0.55, 1.05, 9.00, 0.32, "👉 一起想一想 — AI 是怎么猜的?",
   sz=13, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.38, 9.00, 0.24, "Let's think together — how did AI guess?",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# 4 reflection questions in 2x2 grid
qs = [
    ("👀", "AI 在看什么?",
     "在看你画的 「线条 + 形状」", CYBER),
    ("🤔", "AI 为什么这样猜?",
     "它比较以前看过的很多图", ORANGE),
    ("❓", "AI 真的 「认识」 苹果吗?",
     "不 — 它只在比较模式!", PINK),
    ("📚", "AI 看过几张苹果图?",
     "几千几万张! (Data 训练)", GREEN),
]
card_w = 4.40; gap = 0.30; row_gap = 0.20
total_w = 2 * card_w + gap; start = (10 - total_w) / 2
for i, (em, q, a, cl) in enumerate(qs):
    row = i // 2; col = i % 2
    x = start + col * (card_w + gap)
    y = 1.85 + row * (1.45 + row_gap)
    panel(s, x, y, card_w, 1.45, cl, lw=2.5)
    tb(s, x + 0.15, y + 0.10, 0.55, 0.55, em, sz=28)
    tb(s, x + 0.75, y + 0.15, card_w - 0.85, 0.40, q, sz=13, b=True, c=cl)
    tb(s, x + 0.20, y + 0.80, card_w - 0.40, 0.55, a, sz=11, b=True, c=DARK)

# Key takeaway band
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.60))
band.fill.solid(); band.fill.fore_color.rgb = INK
band.line.color.rgb = STAR; band.line.width = Pt(2)
tb(s, 0.55, 5.02, 9.00, 0.30, "💡 今天的 「啊哈」: AI 会看、会猜、会犯错 — 但它不是真 「懂」!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.30, 9.00, 0.22, "Today's 'aha': AI sees, guesses, makes mistakes — but doesn't really understand",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5-8 分钟讨论:\n• 老师念问题, 学生答\n• K-3: 重点 — 「AI 会看, 会猜, 会犯错」\n• 4-5: 引入 「数据训练」 概念 (Day 3 ML 铺垫)\n\n关键学习点:\n• AI 不是像人一样真 「懂」\n• 它是根据以前看过很多图来猜\n• 这也是 AI 会犯错的原因\n\n准备 Day 3 (Machine Learning):\n• 今天学了 — AI 会看 + 会猜\n• Day 3 — 我们自己 「教」 AI 看 + 猜")

# ============================================================
# 19 · 小侦探三步法
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔍 小侦探三步法  Be a Tiny Detective!", PURPLE)
tb(s, 0.40, 0.85, 9.20, 0.30, "AI 说的对不对? 用这 3 步检查!",
   sz=13, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.40, 1.15, 9.20, 0.24, "Don't believe right away — check with 3 steps!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

steps = [
    ("1️⃣", "🔎", "查来源", "Check Source",
     "这个信息是谁说的? 靠谱吗?", "Who said it? Reliable?", CYBER),
    ("2️⃣", "📚", "多比较", "Compare More",
     "其他书 / 网站也这么说吗?", "Do other sources agree?", ORANGE),
    ("3️⃣", "🧠", "动脑筋", "Use Your Brain",
     "这符合常理吗? 会不会太夸张?", "Does it make sense?", GREEN),
]
card_w = 2.95; gap = 0.12
total = 3 * card_w + 2 * gap; start = (10 - total) / 2
for i, (num, em, cn_t, en_t, cn_d, en_d, cl) in enumerate(steps):
    x = start + i * (card_w + gap)
    panel(s, x, 1.55, card_w, 3.05, cl, lw=3)
    # Number badge
    bg_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.05), Inches(1.62),
                                   Inches(0.55), Inches(0.55))
    bg_circle.fill.solid(); bg_circle.fill.fore_color.rgb = cl
    bg_circle.line.fill.background()
    tb(s, x + 0.05, 1.66, 0.55, 0.45, num, sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Icon
    tb(s, x + 0.70, 1.62, card_w - 0.70, 0.55, em, sz=30)
    # Title
    tb(s, x + 0.15, 2.30, card_w - 0.30, 0.40, cn_t, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.15, 2.70, card_w - 0.30, 0.28, en_t, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Description
    tb(s, x + 0.15, 3.10, card_w - 0.30, 0.80, cn_d, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x + 0.15, 3.95, card_w - 0.30, 0.60, en_d, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# Practice prompt
activity_box(s, 0.40, 4.80, 9.20, 0.75,
             "练一练: 「每天睡 20 小时最健康」 — 用 3 步检查!",
             "Practice: Apply the 3 steps to a wrong AI claim!", color=PURPLE)
n += 1; pn(s, n)
notes(s, "8 分钟:\n• 跟读三步名字 3 次\n• 用 「每天睡 20 小时」 例子实际检查:\n  1. 查来源 → AI 说的, 不是医生\n  2. 多比较 → 医生 / 父母 / 书都说 8-10 小时\n  3. 动脑筋 → 20 小时 = 睡一整天, 不合理!\n• 让学生总结: 哪一步最重要?")

# ============================================================
# 16 · 做聪明的 AI 小主人 (intro to 5 scenarios)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "👑 做聪明的 AI 小主人!  Be a Smart AI Master!", DAY)

# Intro
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.30))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = DAY; intro.line.width = Pt(2.5)
tb(s, 0.55, 1.05, 9.00, 0.40, "🤖 AI 很厉害 — 但你是 「小主人」!",
   sz=15, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.48, 9.00, 0.32, "AI 听你的 — 不是你听 AI 的!",
   sz=13, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.82, 9.00, 0.28, "AI is powerful — but YOU are the smart master!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 4 principles preview
principles = [
    ("🧠", "我会思考", "I Think", CYBER),
    ("🔍", "我会辨别", "I Verify", PURPLE),
    ("🛡️", "我会保护", "I Protect", GREEN),
    ("🎨", "我会创造", "I Create", ORANGE),
]
card_w = 2.20; gap = 0.12
total = 4 * card_w + 3 * gap; start = (10 - total) / 2
for i, (em, cn, en, cl) in enumerate(principles):
    x = start + i * (card_w + gap)
    panel(s, x, 2.45, card_w, 1.95, cl, lw=2.5)
    tb(s, x, 2.60, card_w, 0.75, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 3.40, card_w, 0.42, cn, sz=16, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 3.85, card_w, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)

activity_box(s, 0.40, 4.55, 9.20, 1.00,
             "接下来 5 个情景 — 你来判断: ✅ 聪明还是 ❌ 不聪明?",
             "Next: 5 scenarios — you vote: Smart ✅ or Not ❌?",
             gesture_hint="👏 全班一起投票 + 说 「为什么」",
             color=DAY)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 强调 「小主人」 概念 — 学生是主人, AI 是工具\n• 4 个 「我会」 不用讲太多 — 通过情景学")

# ============================================================
# 17-21 · 5 Smart-AI-User Scenarios
# ============================================================
scenarios = [
    (1, "AI 帮我做数学题, 我直接抄上去!",
     "AI did my math homework, I just copied it!",
     "❌",
     "你没有学到 — AI 帮你, 不是替你! 「思考」 才是你的工作!",
     "You didn't learn — AI helps, doesn't replace YOU!",
     # Teacher notes
     "✅ 答案: ❌ 不聪明\n\n"
     "为什么不聪明:\n"
     "• 抄答案 = 学不到东西 — 下次考试还是不会\n"
     "• AI 帮的是「让你懂」, 不是「替你做」\n"
     "• 学数学 = 学「怎么想」, 不是「拿答案」\n\n"
     "课堂讨论问题 (引导 K-5):\n"
     "• 「如果你抄了, 明天老师问你, 你会答吗?」\n"
     "• 「妈妈帮你拧瓶盖, 你下次还是不会, 对吗?」\n"
     "• 「那 AI 帮做作业像不像妈妈帮拧瓶盖?」\n\n"
     "聪明的做法:\n"
     "• 让 AI 解释步骤 → 你自己再做一遍\n"
     "• 不会的地方让 AI 当小老师, 不当替身\n\n"
     "延伸: 这跟「为什么大家学习」连起来 — 我们是来学「会思考」的, 不是来拿答案的"),

    (2, "AI 帮我解释题目, 我自己做!",
     "AI explained the problem, I solved it myself!",
     "✅",
     "AI 是你的 「老师」, 不是 「替身」 — 你才是学习的主角!",
     "AI is a teacher, not a replacement — YOU are the learner!",
     # Teacher notes
     "✅ 答案: ✅ 聪明小主人!\n\n"
     "为什么聪明:\n"
     "• AI 是你的「24 小时小老师」 — 不懂可以随时问\n"
     "• 你自己做 = 大脑真的在学习\n"
     "• 下次遇到类似题, 你自己就会了!\n\n"
     "课堂讨论问题:\n"
     "• 「这样用 AI 和场景 1 哪里不一样?」 (重点对比!)\n"
     "• 「如果你已经懂了, 还要做练习吗?」 (要 — 加深记忆)\n"
     "• 「AI 解释 + 你自己做 = 谁是主角?」 (你!)\n\n"
     "推荐用法 (告诉学生):\n"
     "• 题目看不懂 → 让 AI 用简单话讲一遍\n"
     "• 不会做 → 让 AI 给一个相似的例题\n"
     "• 做错了 → 让 AI 帮你找哪一步错了\n\n"
     "延伸: 这是 AI 最好的用法 — 像有了无数个家教!"),

    (3, "我把自己的电话 + 地址告诉 AI!",
     "I told AI my phone number and address!",
     "❌",
     "🛡️ 不要把 「小秘密」 告诉 AI! 隐私要保护 — 像保护你的玩具!",
     "Never share personal info with AI! Privacy first!",
     # Teacher notes
     "✅ 答案: ❌ 不聪明 (而且不安全!)\n\n"
     "为什么不安全:\n"
     "• 网上有坏人 — 不知道是不是 AI 在听\n"
     "• 电话/地址 = 你家的「钥匙」 — 不能随便给\n"
     "• 即使是真 AI, 你的信息可能被「记住」\n\n"
     "课堂讨论问题:\n"
     "• 「你会把家钥匙给陌生人吗?」 → 不会!\n"
     "• 「电话/地址跟钥匙像不像?」 → 像!\n"
     "• 「那应该告诉 AI 吗?」 → 不能!\n\n"
     "「小秘密」清单 (反复念, 让学生记住):\n"
     "• 🚫 真实姓名 + 家庭住址\n"
     "• 🚫 学校名字 + 班级\n"
     "• 🚫 爸爸妈妈的电话/密码\n"
     "• 🚫 自己和家人的脸部照片\n\n"
     "记住这句话: 「网络世界, 不透露 = 最好的防护盾!」\n\n"
     "延伸 (高年级): 引入 「数据隐私」 概念 — AI 公司可能收集你的对话"),

    (4, "我让 AI 帮我想作文的点子!",
     "I asked AI to help me brainstorm essay ideas!",
     "✅",
     "AI 是你的 「头脑风暴伙伴」 — 但最后写是你自己!",
     "AI is a brainstorm buddy — but YOU write it!",
     # Teacher notes
     "✅ 答案: ✅ 聪明小主人!\n\n"
     "为什么聪明 (但有条件!):\n"
     "• 「想点子」 ≠ 「写作文」 — 这是关键区别\n"
     "• AI 给 10 个点子 → 你选一个 → 你自己写 → 你的作文!\n"
     "• 创意还是你的, AI 只是「打开思路」\n\n"
     "课堂讨论问题:\n"
     "• 「想点子和写作文哪个更难?」 (因人而异 — 都要练)\n"
     "• 「如果 AI 直接写完, 这还是你的作文吗?」 (不是!)\n"
     "• 「问 AI 「写一篇关于夏天的作文」 vs 「夏天可以写什么主题?」 — 哪个更聪明?」 (后者)\n\n"
     "聪明的问 AI 方法 (示范!):\n"
     "• ❌ 笨: 「帮我写作文」\n"
     "• ✅ 聪明: 「我要写夏天的作文 — 可以从哪些角度入手?」\n"
     "• ✅ 更聪明: 「我已经想到游泳、冰淇淋 — 还有什么有趣的角度?」\n\n"
     "延伸: 这呼应了 reference deck 中的 「回声魔镜」 概念 — 你问得越清楚, AI 答得越棒"),

    (5, "AI 帮我画图, 我再修改 + 创作!",
     "AI made an image, I improved & added my ideas!",
     "✅",
     "这是 「人 + AI 合作」 — 最棒的用法! 创意是你的!",
     "Human + AI = best combo! Your creativity counts!",
     # Teacher notes
     "✅ 答案: ✅ 聪明小主人! (最棒的用法!)\n\n"
     "为什么最棒:\n"
     "• 「修改 + 创作」 = 你在主导, AI 是工具\n"
     "• AI 出图快 (10 秒), 你来添加「人的味道」\n"
     "• 像用画笔/铅笔一样 — 工具是工具, 创作者是你\n\n"
     "课堂讨论问题:\n"
     "• 「直接用 AI 图 vs 修改后用 — 哪个更是你的作品?」\n"
     "• 「画家会用电脑画画吗?」 → 会! (Photoshop, iPad)\n"
     "• 「电脑画 vs 手画 — 都算艺术吗?」 → 都算!\n\n"
     "对比 3 种用法 (让学生选最聪明):\n"
     "• A. 直接拿 AI 图当自己作品 → ❌ 不诚实\n"
     "• B. AI 出图 + 你不改 → ⚠️ 不算「你的」作品\n"
     "• C. AI 出图 + 你大改 / 加新元素 → ✅ 聪明合作\n\n"
     "今天的项目 (Session 3) 就是这种用法 — 学生设计, 老师用 AI 工具出插图, 学生再加创意!\n\n"
     "延伸: 这是 「AI 时代」 的新技能 — 不是抗拒 AI, 是和它合作"),
]
for num, scene_cn, scene_en, verdict, reason_cn, reason_en, teacher_note in scenarios:
    s = ns(prs)
    scenario_slide(s, num, scene_cn, scene_en, verdict, reason_cn, reason_en)
    notes(s, teacher_note)
    n += 1; pn(s, n)

# ============================================================
# 22 · AI 使用公约
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📜 我们的 AI 使用公约!  Our AI Promise!", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "我承诺, 做一个聪明的 AI 小主人:",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 1.18, 9.20, 0.24, "I promise to be a smart AI master!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

vows = [
    ("🧠", "我会思考", "I Think",
     "把 AI 当 「参谋」, 不当 「作业替身」", CYBER),
    ("🔍", "我会辨别", "I Verify",
     "对信息保持好奇, 也保持警惕", PURPLE),
    ("🛡️", "我会保护", "I Protect",
     "守护自己和别人的 「小秘密」", GREEN),
    ("🎨", "我会创造", "I Create",
     "用 AI 帮我实现奇思妙想", ORANGE),
]
card_w = 4.40; gap = 0.30; row_gap = 0.20
total_w = 2 * card_w + gap; start = (10 - total_w) / 2
for i, (em, cn_t, en_t, cn_d, cl) in enumerate(vows):
    row = i // 2; col = i % 2
    x = start + col * (card_w + gap)
    y = 1.55 + row * (1.75 + row_gap)
    panel(s, x, y, card_w, 1.75, cl, lw=2.5)
    tb(s, x + 0.15, y + 0.15, 0.85, 0.85, em, sz=44)
    tb(s, x + 1.05, y + 0.20, 3.20, 0.45, cn_t, sz=20, b=True, c=cl)
    tb(s, x + 1.05, y + 0.65, 3.20, 0.28, en_t, sz=10, c=GRAY)
    tb(s, x + 0.20, y + 1.05, card_w - 0.40, 0.60, cn_d, sz=11, b=True, c=DARK)

# Bottom: signature line
sig = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.45))
sig.fill.solid(); sig.fill.fore_color.rgb = DAY; sig.line.fill.background()
tb(s, 0.55, 5.16, 9.00, 0.32, "✍️ 签名: ________________  ·  日期: ____ /____ /____",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟:\n• 老师念 4 个 「我会」, 全班跟着念\n• 学生可以在 booklet 上签名 (Session 2)\n• 拍集体 「公约照」 — 鼓励仪式感\n• 总结 Session 1, 准备下一节")

# ============================================================
# 23 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2", "📖 下午 2:00–2:45  ·  复习 + 中文词汇 + Mini Booklet", DAY, "📚")
n += 1; pn(s, n)

# ============================================================
# 24 · 早上回顾 (recap)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔁 早上学了什么?  Morning Recap", DAY)
tb(s, 0.40, 0.85, 9.20, 0.30, "想一想 — 我们早上学了什么? 一起喊出来!",
   sz=13, b=True, c=DAY, a=PP_ALIGN.CENTER)

recap = [
    ("🤖", "AI 是什么?", "聪明的小助手 — 能听、看、答问题、解决问题", CYBER),
    ("🌍", "哪些算 AI?", "Siri / 自动驾驶 / ChatGPT / Maps  ·  电灯 ✗", ORANGE),
    ("🔍", "AI 会犯错!", "AI 幻觉 — 用 「三步法」 检查!", PURPLE),
    ("👑", "我是小主人!", "我会思考 / 辨别 / 保护 / 创造", GREEN),
]
for i, (em, q, a, cl) in enumerate(recap):
    y = 1.30 + i * 0.95
    panel(s, 0.50, y, 9.00, 0.80, cl, fill=WHITE, lw=2)
    tb(s, 0.65, y + 0.18, 0.55, 0.50, em, sz=24)
    tb(s, 1.30, y + 0.10, 8.00, 0.32, q, sz=13, b=True, c=cl)
    tb(s, 1.30, y + 0.42, 8.00, 0.32, a, sz=10, b=True, c=DARK)
n += 1; pn(s, n)
notes(s, "3-5 分钟:\n• 快速回顾 — 让学生喊答案\n• 老师念问题, 学生答\n• 给 K-2 时间想, 3-5 抢答")

# ============================================================
# 25-28 · 我会认 (4 vocab cards)
# ============================================================
recognize_words = [
    ("🤖", "人工智能", "rén gōng zhì néng", "AI",
     "AI 可以帮我们写故事!", "AI can help us write stories!",
     "ChatGPT / Siri 截图 + AI 机器人图", DAY),
    ("🤖", "机器人", "jī qì rén", "Robot",
     "这个机器人会扫地。", "This robot can sweep the floor.",
     "扫地机器人 / 工业机器人 / 卡通机器人", CYBER),
    ("💻", "电脑", "diàn nǎo", "Computer",
     "我用电脑学习中文。", "I use a computer to learn Chinese.",
     "笔记本 / 台式电脑", ORANGE),
    ("📚", "学习", "xué xí", "Learn",
     "AI 也在学习 — 像我们!", "AI is also learning — like us!",
     "小朋友在看书 / 上课", GREEN),
]
for em, cn, py, en, ex_cn, ex_en, hint, cl in recognize_words:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)

# ============================================================
# 29-30 · 我会写 (2 writing slides)
# ============================================================
s = vocab_write(prs, CYBER, "电脑", "Computer",
                [("电", "diàn", "5 笔", "上面 「日」 + 下面弯钩"),
                 ("脑", "nǎo", "10 笔", "「月」 旁 + 「凶」 头")])
n += 1; pn(s, n)

s = vocab_write(prs, GREEN, "学习", "Learn",
                [("学", "xué", "8 笔", "上 「⺍」 中 「冖」 下 「子」"),
                 ("习", "xí", "3 笔", "像鸟翅膀 — 反复练习")])
n += 1; pn(s, n)

# ============================================================
# SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3", "🎨 下午 3:00–4:30  ·  AI 故事书 + 创意设计", DAY, "🚀")
n += 1; pn(s, n)

# Complete today's booklet — before project starts
s = booklet_slide(prs, day_num=1, day_topic_cn="认识 AI · 做 聪明 的 AI 主人", day_color=DAY)
n += 1; pn(s, n)

# ============================================================
# 主项目 · 分组用 AI 写故事 + AI 配图  (4 slides)
# ============================================================

# ----- Slide A · Project intro -----
s = ns(prs); bg(s, CREAM, prs); hb(s, "📖 主项目 · 分组用 AI 写故事 + AI 配图!", DAY)

# Intro card
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.10))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = DAY; intro.line.width = Pt(2.5)
tb(s, 0.55, 1.05, 9.00, 0.40, "🎯 小组一起写一个故事 → 老师用 AI 给每一页配图!",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.50, 9.00, 0.28, "Groups write a story together → teacher uses AI to illustrate each page",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.78, 9.00, 0.26, "👥 2-4 组 · 每组 4-5 人  ·  ⏱️ 60 min 写 + 配图 + 20 min 展示",
   sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 5 game rules (boxes)
rules = [
    ("1️⃣", "👥", "分小组", "每组 4-5 人  ·  保证一台电脑 + 一个会打中文的同学", CYBER),
    ("2️⃣", "✏️", "选主题", "太空? 动物? 神奇校车? — 一起决定", ORANGE),
    ("3️⃣", "📝", "每人写一句", "组长打字 + 加上每个人的名字!", GREEN),
    ("4️⃣", "🎨", "AI 配图", "老师把故事输入 ChatGPT / DALL-E → 每页一张图", PINK),
    ("5️⃣", "🎤", "上台展示", "每人讲解自己写的那一页!", PURPLE),
]
card_w = 1.80; gap = 0.05
total_w = 5 * card_w + 4 * gap; start = (10 - total_w) / 2
for i, (num, em, cn_t, cn_d, cl) in enumerate(rules):
    x = start + i * (card_w + gap)
    panel(s, x, 2.30, card_w, 2.45, cl, lw=2)
    tb(s, x, 2.40, card_w, 0.30, num, sz=14, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, card_w, 0.55, em, sz=28, a=PP_ALIGN.CENTER)
    tb(s, x + 0.05, 3.30, card_w - 0.10, 0.35, cn_t, sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.10, 3.70, card_w - 0.20, 1.00, cn_d, sz=8, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Tool bar
tool = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.62))
tool.fill.solid(); tool.fill.fore_color.rgb = DAY; tool.line.fill.background()
tb(s, 0.55, 5.02, 9.00, 0.30, "🛠️ 工具: ChatGPT (DALL·E) · DeepSeek · Bing Image Creator (免费!)",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.32, 9.00, 0.22, "Tools: ChatGPT, DeepSeek, Bing Image (free) — all generate illustrations",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟介绍 + 35 分钟写作 + 20 分钟 AI 配图 + 20 分钟展示\n\n准备:\n• 每组一台电脑 / iPad (能登录 Google Doc + AI 工具)\n• 投影连接 — 大家一起看 AI 出图\n• 中文打字法 — 保证每组至少一个会打字的学生 (或老师代打)\n\n分组建议:\n• 4-5 人一组 (人多更热闹, 但写作时间也更长)\n• 混合年龄 — 大孩子帮小孩子写句子\n• 让会中文打字的当 「组长」\n\n时间安排:\n• 5 min — 介绍规则\n• 10 min — 选主题 + 讨论故事大纲\n• 25 min — 每人写一句 (轮流)\n• 20 min — 老师配图 (一边写一边配)\n• 20 min — 上台展示")

# ----- Slide B · Story-writing 5 steps -----
s = ns(prs); bg(s, CREAM, prs); hb(s, "📖 故事接龙 · 怎么写  How to Write Together", PURPLE)
tb(s, 0.40, 0.85, 9.20, 0.30, "每人一句, 接着上一个人 — 一个连贯的故事!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.40, 1.15, 9.20, 0.24, "Each student adds one sentence — connecting to the last",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# 5-page structure (one student per page)
pages = [
    ("Page 1", "🌟", "开头", "Beginning",
     "「有一个 ___ 的 ___」", CYBER),
    ("Page 2", "🎈", "发生了", "Something happens",
     "「突然, 它 ___」", ORANGE),
    ("Page 3", "🚀", "去了哪里", "Goes somewhere",
     "「它去了 ___」", GREEN),
    ("Page 4", "💫", "遇到了", "Meets someone",
     "「它遇到了 ___」", PINK),
    ("Page 5", "🌈", "结尾", "Happy ending",
     "「最后 ___ 开心地 ___」", PURPLE),
]
card_w = 1.78; gap = 0.10
total_w = 5 * card_w + 4 * gap; start = (10 - total_w) / 2
for i, (label, em, cn_t, en_t, frame, cl) in enumerate(pages):
    x = start + i * (card_w + gap)
    panel(s, x, 1.55, card_w, 3.30, cl, lw=2.5)
    # Page label header
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(1.55), Inches(card_w), Inches(0.40))
    head.fill.solid(); head.fill.fore_color.rgb = cl; head.line.fill.background()
    tb(s, x, 1.60, card_w, 0.32, label, sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Icon
    tb(s, x, 2.05, card_w, 0.65, em, sz=30, a=PP_ALIGN.CENTER)
    # Title
    tb(s, x + 0.05, 2.75, card_w - 0.10, 0.35, cn_t, sz=13, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x + 0.05, 3.10, card_w - 0.10, 0.25, en_t, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Sentence frame
    fr_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x + 0.10), Inches(3.50), Inches(card_w - 0.20), Inches(1.20))
    fr_box.fill.solid(); fr_box.fill.fore_color.rgb = WARM
    fr_box.line.color.rgb = cl; fr_box.line.width = Pt(1)
    tb(s, x + 0.15, 3.65, card_w - 0.30, 0.90, frame, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Activity callout
activity_box(s, 0.40, 4.95, 9.20, 0.62,
             "💡 不会写? 用「Roll-a-Story」骰子玩法 — 每人骰一个: 主角 / 地点 / 难题!",
             "Stuck? Try Roll-a-Story dice — each rolls: character / place / problem",
             color=PURPLE)
n += 1; pn(s, n)
notes(s, "25 分钟写作:\n\n方法 1 — 接龙:\n• 第一个学生写「开头」, 后面每人接一句\n• 老师走动帮助 (尤其 K-3)\n• 句型框架可贴在桌上\n\n方法 2 — Roll-a-Story:\n• 准备 3 个 「骰子」 (或抽签):\n  - 骰子 1: 主角 (小猫/外星人/机器人/小朋友...)\n  - 骰子 2: 地点 (月球/森林/学校/海底...)\n  - 骰子 3: 难题 (迷路了/没有朋友/找东西...)\n• 每组骰 → 故事自动有了起点\n\n分层:\n• K-2: 老师把句型贴在桌上, 学生填空\n• 3-5: 自由写, 鼓励加细节\n\n确保:\n• 每个学生都写了至少一句\n• 组长把每人名字打在自己的句子后面\n• 故事整体要 「能讲通」 — 老师可以帮忙调整")

# ----- Slide C · Story example -----
s = ns(prs); bg(s, CREAM, prs); hb(s, "📚 故事示例 · 小外星人变小鸟  Example Story", CYBER)
tb(s, 0.40, 0.85, 9.20, 0.30, "5 个学生一起写的故事 — 你也可以这样写!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 5 sentences from reference deck, one per student
sentences = [
    ("Eva", "🛸", "有一个快乐的小外星人。它每天在星球上玩, 非常开心。一天, 它觉得无聊了, 想去别的地方看看。", CYBER),
    ("John", "🐦", "突然, 它变成了一只蓝色的小鸟! 小鸟拍拍翅膀, 飞了起来。它飞过高高的山, 飞过大大的海。", ORANGE),
    ("Alice", "🌙", "它一边飞一边想: 「我要去月亮看看!」小鸟用力飞啊飞, 终于到了月亮上。", GREEN),
    ("Max", "🐰", "月亮上很安静。地上亮亮的, 软软的。小鸟看到有几只可爱的月亮兔在跳舞。", PINK),
    ("Kyle", "🌟", "小鸟跟月亮兔一起玩, 一起跳, 一起笑。它觉得非常开心。玩累了, 小鸟躺在月亮上, 看着星星, 慢慢睡着了。", PURPLE),
]
for i, (name, em, line, cl) in enumerate(sentences):
    y = 1.30 + i * 0.72
    panel(s, 0.40, y, 9.20, 0.62, cl, lw=2)
    # Page badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(y + 0.10),
                              Inches(0.45), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = cl; badge.line.fill.background()
    tb(s, 0.55, y + 0.16, 0.45, 0.30, str(i+1), sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Emoji
    tb(s, 1.10, y + 0.10, 0.55, 0.45, em, sz=20)
    # Sentence
    tb(s, 1.70, y + 0.08, 7.10, 0.50, line, sz=10, b=True, c=DARK)
    # Author
    tb(s, 8.85, y + 0.20, 0.65, 0.30, f"-{name}", sz=10, b=True, c=cl, a=PP_ALIGN.RIGHT)

# Bottom tip
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.62))
tip.fill.solid(); tip.fill.fore_color.rgb = INK; tip.line.color.rgb = STAR; tip.line.width = Pt(2)
tb(s, 0.55, 5.02, 9.00, 0.28, "✨ 看! 每个人写一句 — 加起来就是一个完整的故事!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.32, 9.00, 0.20, "See! Each student writes one sentence — together = one whole story!",
   sz=8, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "用这个示例:\n• 老师念一遍整个故事 (或让学生轮流读)\n• 让学生注意 — 句子怎么 「连得起来」\n• 强调: 每人只写一句, 但加起来变成一个完整的故事!\n\n讨论问题:\n• 「Eva 写的是什么? 是开头还是结尾?」\n• 「John 怎么接 Eva 的句子的?」\n• 「你最喜欢哪一句? 为什么?」\n\n灵感主题 (替换 「外星人」):\n• 一只迷路的小狗\n• 一个会魔法的小书包\n• 一个住在树洞里的精灵\n• 一个发明家小朋友\n• 一个生病的小恐龙")

# ----- Slide D · Teacher AI prompt workflow -----
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎨 老师工作流 · 怎么让 AI 画好图  Teacher's AI Workflow", ORANGE)

# Top intro
intro = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.70))
intro.fill.solid(); intro.fill.fore_color.rgb = WARM
intro.line.color.rgb = ORANGE; intro.line.width = Pt(2)
tb(s, 0.55, 1.02, 9.00, 0.30, "🎯 学生写完一句 → 老师立刻让 AI 画图 → 当场展示!",
   sz=13, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.35, 9.00, 0.24, "Student writes a sentence → teacher prompts AI → live illustration!",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: 3 prompt rules
panel(s, 0.40, 1.80, 4.55, 3.20, ORANGE)
panel_head(s, 0.40, 1.80, 4.55, ORANGE, "✏️ 写好 Prompt 的 3 个秘诀", sz=12)
rules = [
    ("1️⃣", "明确生成什么", "「请画一张图...」"),
    ("2️⃣", "描述细节", "颜色 + 动作 + 表情 + 背景"),
    ("3️⃣", "加风格", "卡通 / 水彩 / 吉卜力风格"),
]
for i, (num, cn_t, cn_d) in enumerate(rules):
    y = 2.40 + i * 0.82
    tb(s, 0.55, y, 0.50, 0.35, num, sz=14, b=True, c=ORANGE)
    tb(s, 1.10, y, 3.70, 0.35, cn_t, sz=12, b=True, c=DARK)
    tb(s, 1.10, y + 0.38, 3.70, 0.40, cn_d, sz=10, c=GRAY)

# RIGHT: Good vs Bad prompt examples
panel(s, 5.05, 1.80, 4.55, 3.20, PURPLE)
panel_head(s, 5.05, 1.80, 4.55, PURPLE, "💡 例子 · 笨 vs 聪明", sz=12)
# Bad
bad_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(5.20), Inches(2.35), Inches(4.30), Inches(0.95))
bad_box.fill.solid(); bad_box.fill.fore_color.rgb = WHITE
bad_box.line.color.rgb = PINK; bad_box.line.width = Pt(1.5)
tb(s, 5.30, 2.42, 4.10, 0.25, "❌ 笨:", sz=10, b=True, c=PINK)
tb(s, 5.30, 2.65, 4.10, 0.55, "「画一只猫」", sz=12, b=True, c=DARK)
# Good
good_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(5.20), Inches(3.42), Inches(4.30), Inches(1.45))
good_box.fill.solid(); good_box.fill.fore_color.rgb = WHITE
good_box.line.color.rgb = GREEN; good_box.line.width = Pt(1.5)
tb(s, 5.30, 3.50, 4.10, 0.25, "✅ 聪明:", sz=10, b=True, c=GREEN)
tb(s, 5.30, 3.75, 4.10, 1.10, "「请画一只白色的小猫, 在窗台上睡觉, 阳光照在身上, 卡通风格」",
   sz=10, b=True, c=DARK)

# Tools bar
tools = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.45))
tools.fill.solid(); tools.fill.fore_color.rgb = ORANGE; tools.line.fill.background()
tb(s, 0.55, 5.17, 9.00, 0.30, "🛠️ 工具: ChatGPT (4o) · DALL·E · Bing Image Creator · Adobe Firefly",
   sz=11, b=True, c=WHITE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "老师 prompt 写作流程 — 现场演示:\n\n• 学生写: 「小鸟飞过高高的山」\n• 老师在 ChatGPT 输入:\n  「请画一只蓝色的小鸟, 在高山上空飞过, 翅膀张开, 卡通风格, 适合小朋友」\n• AI 30 秒出图 → 投影展示!\n\nPrompt 模板:\n「请生成 [图片/卡通图]: [主角描述] 在 [地点] [动作]. 风格: [卡通/水彩/吉卜力].」\n\n图像生成 prompt 三要素 (源自 reference PPT):\n\n1. 明确生成类型\n   • 「请生成一张图」 / 「画一只」\n   • 改细节: 「把这只小猫变成蓝色」\n\n2. 核心特征描述\n   • 外貌 (毛色、花纹、体型)\n   • 动作 (跑、跳、睡觉、扑玩具)\n   • 表情 (好奇、困倦、惊讶)\n   • 环境 (沙发、花园、窗台)\n\n3. 加风格\n   • 写实 / 卡通 / 水彩 / 吉卜力\n   • 「皮克斯风格 3D 渲染」 (高级用法)\n\n推荐工具:\n• ChatGPT (GPT-4o 自带图像) — 最方便\n• Bing Image Creator (免费, bing.com/create)\n• DALL·E (openai.com)\n• Adobe Firefly (firefly.adobe.com)\n\n小诀窍:\n• 让学生看 AI 出图过程 — 学生很惊喜\n• 如果 AI 画得不像, 当场修改 prompt 再试一次\n• 这本身就是教学 — 学生看到 「怎么让 AI 听懂」")

# ============================================================
# 37 · Share + Close
# ============================================================
s = share_close(prs, DAY,
    frames_cn=["「我的 AI 叫 ___, 它帮 ___」",
               "「聪明主人怎么用它? ___」"],
    frames_en="My AI is called ___, it helps ___ · A smart master uses it by ___",
    next_day_cn="Day 2 · 3D 打印 — 从想象到真实!",
    next_day_en="Day 2 · 3D Printing — From idea to real object!",
    next_emoji="🖨️")
n += 1; pn(s, n)
notes(s, "15-20 分钟分享:\n• 每组 1 人上台介绍 (1-2 分钟)\n• 全班鼓掌\n• 老师拍照留念 — 集体照 + 作品照\n• 提醒明天主题: 3D 打印\n• 让学生把故事书带回家给家长看")

# ============================================================
out = os.path.join(os.path.dirname(__file__), "day1_ai.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
