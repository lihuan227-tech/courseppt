#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
玩转 创新 科技 · Day 3: 机器 学习 (Machine Learning)
3-session classroom deck for K-5 Chinese immersion summer camp.

探究 问题: 电脑 怎么 学 习?

我 会 认: 机器 / 分类 / 数据 / 训练
我 会 写: 分类 / 训练

Structure (matches Day 1 / Day 2):
  Session 1 (11:00–11:45) — 认识 机器 学习  · What is ML?
  Session 2 (2:00–2:45)   — 复习 + 中文 词汇  · 我 会 认 / 我 会 写
  Session 3 (3:00–4:30)   — 训练 真 AI + 我 是 AI 训练师  · Project
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = make_presentation()
DAY = ML_GREEN
n = 0


def arrow(s, x, y, w=0.30, h=0.30, color=DAY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ============================================================
# 1 · COVER
# ============================================================
cover(prs, 3, "Machine Learning", "机器 学习 · 电脑 怎么 学 习?",
      "🧠 📊 🐱 🤖 ✨", DAY,
      "电脑 怎么 学 习?",
      "How do computers learn?")
n += 1; pn(prs.slides[-1], n)
notes(prs.slides[-1], "Day 3 · 核心 概念: 看 例子 → 找 规律 → 猜 答案\n• Session 1: 认识 机器 学习 + 真 实 例子\n• Session 2: 中文 词汇 (我 会 认 / 我 会 写)\n• Session 3: 用 Teachable Machine 训练 真 AI + 项目")


# ============================================================
# 2 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1", "🌅 上午 11:00–11:45  ·  机器 学习 实验 室 · ML Lab",
        DAY, "🔬"); n += 1; pn(s, n)


# ============================================================
# 3 · LEARNING GOALS
# ============================================================
s = learning_goals(prs, DAY, [
    ("1️⃣", "通过 真 实 实验 体验 机器 学习",
     "Experience ML through a real experiment", CYBER),
    ("2️⃣", "知道 「数据 / 分类 / 特征 / 训练」 是 什么",
     "Learn: data, classify, features, training", ORANGE),
    ("3️⃣", "发现 AI 为什么 有 时 会 猜 错",
     "Discover why AI sometimes makes mistakes", PINK),
    ("4️⃣", "学 当 一 个 聪明 的 AI 训练 师",
     "Become a thoughtful AI trainer", PURPLE),
])
n += 1; pn(s, n)


# ============================================================
# 4 · LAB OPENING TITLE — 机器 学习 实验 室
# ============================================================
s = ns(prs); bg(s, INK, prs)
# Lab decoration: corner sparkles + beakers
for x, y in [(0.5, 0.5), (9.0, 0.5), (0.5, 4.8), (9.0, 4.8)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()
# Top banner
tb(s, 0.3, 0.45, 9.4, 0.40, "🔬 Session 1 · Machine Learning Lab",
   sz=16, b=True, c=NEON, a=PP_ALIGN.CENTER)
# Big title
title_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.05), Inches(8.8), Inches(2.20))
title_box.fill.solid(); title_box.fill.fore_color.rgb = DAY
title_box.line.color.rgb = STAR; title_box.line.width = Pt(4)
tb(s, 0.8, 1.25, 8.4, 0.85, "机器 学习 实验 室",
   sz=46, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.15, 8.4, 0.50, "我 是 AI 训练 师",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.75, 8.4, 0.40, "Machine Learning Lab · I'm an AI Trainer!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)
# Lab icons row
tb(s, 0.3, 3.65, 9.4, 1.10, "🔬   🧪   🤖   💻   ✨",
   sz=58, a=PP_ALIGN.CENTER)
# Bottom hype
tb(s, 0.3, 5.00, 9.4, 0.40, "今天 我们 一起 做 真 实 验, 训练 真 AI!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "30 秒 hype slide.\n• 让 学生 大 声 喊 「我 是 AI 训练 师!」\n• 今天 重点: 不 是 听 讲 — 是 做 实验\n• 老师 准备: 笔记本 + 投影 + 摄像头 + Teachable Machine 已 打开")


# ============================================================
# 5 · HOOK STORY — Baby Learning
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🍼 小 baby 怎么 学 习? · How Does a Baby Learn?", DAY)

# Top story panel
story = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(2.05))
story.fill.solid(); story.fill.fore_color.rgb = WARM
story.line.color.rgb = DAY; story.line.width = Pt(2.5)
# Baby + parent emojis
tb(s, 0.55, 1.05, 1.20, 1.20, "👶", sz=64, a=PP_ALIGN.CENTER)
# Parent teaching dialogue
tb(s, 1.80, 1.05, 7.65, 0.35, "爸爸 妈妈 教 小 baby:",
   sz=13, b=True, c=DAY)
parent_lines = [
    ("🐱", "「这 是 猫。」"),
    ("🐶", "「这 是 狗。」"),
    ("🍎", "「这 是 苹果。」"),
    ("🍌", "「这 是 香蕉。」"),
]
for i, (em, line) in enumerate(parent_lines):
    col = i % 2; row = i // 2
    x = 1.85 + col * 3.85
    y = 1.40 + row * 0.62
    tb(s, x, y, 0.40, 0.40, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, x+0.45, y+0.05, 3.30, 0.40, line, sz=14, b=True, c=DARK)
tb(s, 1.80, 2.65, 7.65, 0.30, "✨ 看 很 多 次 — baby 学 会 了!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)

# Big question + 4 idea bubbles
q = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.20), Inches(9.20), Inches(0.55))
q.fill.solid(); q.fill.fore_color.rgb = DAY; q.line.fill.background()
tb(s, 0.55, 3.30, 9.0, 0.40, "🤔 小 baby 是 怎么 学 会 的?  How did baby learn?",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)

ideas = [
    ("👀", "看 很 多 次"),
    ("👨‍👩‍👧", "别 人 教"),
    ("🧠", "记 住 了"),
    ("🔍", "找 规律"),
]
idea_w = 2.10; gap = 0.15
total = 4*idea_w + 3*gap; start = (10 - total)/2
for i, (em, txt) in enumerate(ideas):
    x = start + i*(idea_w + gap)
    panel(s, x, 3.95, idea_w, 0.85, ORANGE, fill=WHITE, lw=2.5)
    tb(s, x, 4.02, idea_w, 0.40, em, sz=22, a=PP_ALIGN.CENTER)
    tb(s, x, 4.42, idea_w, 0.32, txt, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Bottom takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.45))
tk.fill.solid(); tk.fill.fore_color.rgb = DAY
tk.line.color.rgb = STAR; tk.line.width = Pt(2)
tb(s, 0.55, 5.00, 9.0, 0.35, "💡 电脑 也 是 这 样 学 的!  Computers learn the same way!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 故事 引入:\n• 老师 用 讲 故事 的 语气 慢 慢 说\n• 问: 小 baby 是 怎么 学 会 的?\n• 让 学生 想 + 答 — 别 急 着 给 答案\n• 引出: 电脑 也 是 看 很 多 例子 学 — 叫 「机器 学习」")


# ============================================================
# 6 · OVERVIEW — ML in 6 Steps  (preview before deep dive)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "💡 机器 学习 6 步 · ML in 6 Steps", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "今天 我们 一步 一步 看 — 学 完 这 6 步, AI 就 会 自己 猜!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 6 mini cards: 2 rows of 3
overview_steps = [
    ("1️⃣", "📷", "收集 例子", "Collect", CYBER),
    ("2️⃣", "🏷️", "贴 标签", "Label", ORANGE),
    ("3️⃣", "🧠", "AI 学 习", "Study", PURPLE),
    ("4️⃣", "🔍", "找 规律", "Patterns", PINK),
    ("5️⃣", "🧪", "测 试", "Test", DAY),
    ("6️⃣", "✨", "做 判 断", "Predict", GREEN),
]
ow = 2.85; ogap_x = 0.20; ogap_y = 0.25; oh = 1.65
ostart_x = (10 - 3*ow - 2*ogap_x)/2
for i, (num, em, cn, en, cl) in enumerate(overview_steps):
    row = i // 3; col = i % 3
    x = ostart_x + col*(ow + ogap_x)
    y = 1.30 + row*(oh + ogap_y)
    panel(s, x, y, ow, oh, cl, fill=WHITE, lw=2.5)
    # Number badge top-left
    tb(s, x+0.10, y+0.08, 0.55, 0.45, num, sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
    # Big emoji centered
    tb(s, x+0.60, y+0.10, ow-0.70, 0.75, em, sz=42, a=PP_ALIGN.CENTER)
    # Chinese name
    tb(s, x, y+0.90, ow, 0.40, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    # English subtitle
    tb(s, x, y+1.28, ow, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Arrow to next (same row)
    if col < 2:
        arrow(s, x + ow + 0.02, y + oh/2 - 0.13, w=0.16, h=0.26, color=DAY)
# Down arrow between rows (between cards 3 and 4) — only visible if user looks
# (skipping to avoid visual clutter)

# Bottom takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.15), Inches(9.20), Inches(0.30))
tk.fill.solid(); tk.fill.fore_color.rgb = DAY; tk.line.fill.background()
tb(s, 0.55, 5.18, 9.0, 0.25, "🎬 现在 我们 一 步 一 步 看 详细!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 让 学生 一起 数 「1, 2, 3, 4, 5, 6 步!」\n• 介绍: 「下 面 我们 一 张 slide 看 一 步」\n• 这 是 准备 让 学生 知 道 接 下 来 6 张 slide 的 结 构")


# ============================================================
# 7 · STEP 1 — 收集 例子 (Collect Examples)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📷 Step 1 · 收集 例子 · Collect Examples", CYBER)

# LEFT: visual — grid of cameras/photos
panel(s, 0.40, 0.95, 4.55, 4.10, CYBER, fill=WHITE, lw=3)
# Big number
tb(s, 0.55, 1.05, 4.30, 0.50, "1️⃣", sz=28, b=True, c=CYBER, a=PP_ALIGN.LEFT)
# Big camera grid (10 cameras = "lots of photos")
tb(s, 0.50, 1.65, 4.40, 1.20, "📷 📷 📷 📷 📷",
   sz=46, a=PP_ALIGN.CENTER)
tb(s, 0.50, 2.85, 4.40, 1.20, "📷 📷 📷 📷 📷",
   sz=46, a=PP_ALIGN.CENTER)
tb(s, 0.50, 4.20, 4.40, 0.40, "100 张? 1000 张? 越 多 越 好!",
   sz=12, b=True, c=CYBER, a=PP_ALIGN.CENTER)
tb(s, 0.50, 4.65, 4.40, 0.30, "More photos = better!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# RIGHT: explanation
panel(s, 5.15, 0.95, 4.45, 4.10, WARM, fill=WARM, lw=2)
tb(s, 5.30, 1.10, 4.15, 0.45, "我们 给 AI 很 多 例子!",
   sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.30, 1.58, 4.15, 0.32, "Give AI lots of examples",
   sz=11, c=GRAY, a=PP_ALIGN.LEFT)

# Example panel
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(2.10), Inches(4.15), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
ex.line.color.rgb = CYBER; ex.line.width = Pt(2)
tb(s, 5.40, 2.20, 3.95, 0.30, "🐱 例子:", sz=11, b=True, c=CYBER)
tb(s, 5.40, 2.50, 3.95, 0.35, "拍 100 张 猫 的 照片",
   sz=13, b=True, c=DARK)
tb(s, 5.40, 2.88, 3.95, 0.28, "Take 100 cat photos",
   sz=9, c=GRAY)

# Kid analogy callout
ka = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(3.40), Inches(4.15), Inches(1.50))
ka.fill.solid(); ka.fill.fore_color.rgb = DAY
ka.line.color.rgb = STAR; ka.line.width = Pt(2.5)
tb(s, 5.40, 3.50, 3.95, 0.32, "💡 就 像 ...",
   sz=12, b=True, c=STAR)
tb(s, 5.40, 3.85, 3.95, 0.55, "老师 给 你 100 道 题 练 习!",
   sz=15, b=True, c=WHITE)
tb(s, 5.40, 4.45, 3.95, 0.32, "Like teacher giving you 100 problems to practice!",
   sz=10, c=WARM)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 强调 「越 多 越 好」 — 100 张 比 10 张 强\n• 问 学生: 你 学 拼音 看 了 多 少 次? — 看 很 多 次 才 会!\n• 类比 学校 学 习 让 K-5 易 懂")


# ============================================================
# 8 · STEP 2 — 贴 标签 (Add Labels)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🏷️ Step 2 · 贴 标签 · Add Labels", ORANGE)

# LEFT: photo placeholder + label sticker visual
panel(s, 0.40, 0.95, 4.55, 4.10, ORANGE, fill=WHITE, lw=3)
tb(s, 0.55, 1.05, 4.30, 0.50, "2️⃣", sz=28, b=True, c=ORANGE, a=PP_ALIGN.LEFT)

# Photo placeholder
ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.65), Inches(2.60), Inches(2.00))
ph.fill.solid(); ph.fill.fore_color.rgb = IMGBG
ph.line.color.rgb = ORANGE; ph.line.width = Pt(2)
tb(s, 0.85, 2.30, 2.60, 0.80, "🐱", sz=72, a=PP_ALIGN.CENTER)

# Big sticker label (tilted look using rectangle)
sticker = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.20), Inches(2.40), Inches(1.55), Inches(0.85))
sticker.fill.solid(); sticker.fill.fore_color.rgb = STAR
sticker.line.color.rgb = ORANGE; sticker.line.width = Pt(3)
tb(s, 3.20, 2.45, 1.55, 0.40, "🏷️ 标签", sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 3.20, 2.78, 1.55, 0.45, "「猫」", sz=22, b=True, c=ORANGE, a=PP_ALIGN.CENTER)

# Pair examples below
tb(s, 0.55, 3.95, 4.30, 0.35, "🐱 → 「猫」    🐶 → 「狗」    🍎 → 「苹果」",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.35, 4.30, 0.30, "每 张 都 要 告诉 AI 名字!",
   sz=11, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.65, 4.30, 0.28, "Every photo needs a name (label)",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# RIGHT: explanation
panel(s, 5.15, 0.95, 4.45, 4.10, WARM, fill=WARM, lw=2)
tb(s, 5.30, 1.10, 4.15, 0.45, "告诉 AI 这 是 什么!",
   sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.30, 1.58, 4.15, 0.32, "Tell AI what each one is",
   sz=11, c=GRAY, a=PP_ALIGN.LEFT)

# Explanation box
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(2.10), Inches(4.15), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
ex.line.color.rgb = ORANGE; ex.line.width = Pt(2)
tb(s, 5.40, 2.20, 3.95, 0.30, "📋 重要!", sz=11, b=True, c=ORANGE)
tb(s, 5.40, 2.50, 3.95, 0.35, "如果 不 告诉 AI 名字, AI 就 不 知道!",
   sz=12, b=True, c=DARK)
tb(s, 5.40, 2.90, 3.95, 0.28, "Without labels, AI has no clue",
   sz=9, c=GRAY)

# Kid analogy
ka = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(3.40), Inches(4.15), Inches(1.50))
ka.fill.solid(); ka.fill.fore_color.rgb = DAY
ka.line.color.rgb = STAR; ka.line.width = Pt(2.5)
tb(s, 5.40, 3.50, 3.95, 0.32, "💡 就 像 ...",
   sz=12, b=True, c=STAR)
tb(s, 5.40, 3.85, 3.95, 0.55, "给 玩具 写 上 名字 牌!",
   sz=15, b=True, c=WHITE)
tb(s, 5.40, 4.45, 3.95, 0.32, "Like putting name tags on your toys!",
   sz=10, c=WARM)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 演示: 拿 一 张 图 + 一 张 「标签」 卡片 — 物理 演示\n• 问: 如果 不 给 AI 标签, AI 怎么 知道 是 什么?\n• 引出 「数据 标签」 是 AI 的 「答案」")


# ============================================================
# 9 · STEP 3 — AI 学 习 (AI Studies)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🧠 Step 3 · AI 学 习 · AI Studies", PURPLE)

# LEFT: AI brain + photos flowing in
panel(s, 0.40, 0.95, 4.55, 4.10, PURPLE, fill=WHITE, lw=3)
tb(s, 0.55, 1.05, 4.30, 0.50, "3️⃣", sz=28, b=True, c=PURPLE, a=PP_ALIGN.LEFT)

# Photos flowing in (left side)
tb(s, 0.55, 1.85, 1.30, 0.50, "📷", sz=28, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.45, 1.30, 0.50, "📷", sz=28, a=PP_ALIGN.CENTER)
tb(s, 0.55, 3.05, 1.30, 0.50, "📷", sz=28, a=PP_ALIGN.CENTER)
tb(s, 0.55, 3.65, 1.30, 0.50, "📷", sz=28, a=PP_ALIGN.CENTER)

# Arrows
for ay in [2.05, 2.65, 3.25, 3.85]:
    arrow(s, 1.85, ay, w=0.45, h=0.25, color=PURPLE)

# AI brain (right side of left panel)
brain = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.50), Inches(1.75), Inches(2.20), Inches(2.20))
brain.fill.solid(); brain.fill.fore_color.rgb = WARM
brain.line.color.rgb = PURPLE; brain.line.width = Pt(3)
tb(s, 2.50, 2.10, 2.20, 1.00, "🤖🧠", sz=64, a=PP_ALIGN.CENTER)
tb(s, 2.50, 3.10, 2.20, 0.40, "AI", sz=18, b=True, c=PURPLE, a=PP_ALIGN.CENTER)

tb(s, 0.55, 4.20, 4.30, 0.35, "AI 一 张 一 张 看 ...",
   sz=13, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.55, 4.30, 0.30, "AI looks at them one by one",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# RIGHT: explanation
panel(s, 5.15, 0.95, 4.45, 4.10, WARM, fill=WARM, lw=2)
tb(s, 5.30, 1.10, 4.15, 0.45, "AI 慢 慢 「看」 例子!",
   sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.30, 1.58, 4.15, 0.32, "AI looks at each example",
   sz=11, c=GRAY, a=PP_ALIGN.LEFT)

# Time visual
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(2.10), Inches(4.15), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
ex.line.color.rgb = PURPLE; ex.line.width = Pt(2)
tb(s, 5.40, 2.20, 3.95, 0.30, "⏱️ 慢 一 点!", sz=11, b=True, c=PURPLE)
tb(s, 5.40, 2.50, 3.95, 0.35, "100 张 → 1000 张 → 10000 张!",
   sz=12, b=True, c=DARK)
tb(s, 5.40, 2.88, 3.95, 0.28, "More examples = more practice",
   sz=9, c=GRAY)

# Kid analogy
ka = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(3.40), Inches(4.15), Inches(1.50))
ka.fill.solid(); ka.fill.fore_color.rgb = DAY
ka.line.color.rgb = STAR; ka.line.width = Pt(2.5)
tb(s, 5.40, 3.50, 3.95, 0.32, "💡 就 像 ...",
   sz=12, b=True, c=STAR)
tb(s, 5.40, 3.85, 3.95, 0.55, "你 学 认 字 — 看 多 了 就 会!",
   sz=15, b=True, c=WHITE)
tb(s, 5.40, 4.45, 3.95, 0.32, "Like learning to read — practice makes perfect!",
   sz=10, c=WARM)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 强调: 不 是 一 下子 学 会, 是 慢 慢 看\n• 演示: 老师 假装 自己 是 AI — 拿 起 一 张 一 张 「记」 \n• 引出: 这 个 「学」 的 过程 在 计算 机 里 叫 「training 训练」")


# ============================================================
# 10 · STEP 4 — 找 规律 (Find Patterns)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔍 Step 4 · 找 规律 · Find Patterns", PINK)

# LEFT: AI with magnifying glass + cat features highlighted
panel(s, 0.40, 0.95, 4.55, 4.10, PINK, fill=WHITE, lw=3)
tb(s, 0.55, 1.05, 4.30, 0.50, "4️⃣", sz=28, b=True, c=PINK, a=PP_ALIGN.LEFT)

# Big cat in spotlight
tb(s, 0.55, 1.60, 4.30, 1.40, "🔍   🐱",
   sz=72, a=PP_ALIGN.CENTER)

# Feature checklist below
features_check = [
    ("✅", "尖 耳朵"),
    ("✅", "胡须"),
    ("✅", "尾巴 弯"),
]
for i, (check, feat) in enumerate(features_check):
    y = 3.20 + i*0.45
    tb(s, 1.00, y, 0.40, 0.40, check, sz=18, a=PP_ALIGN.LEFT)
    tb(s, 1.45, y+0.05, 3.30, 0.35, feat, sz=14, b=True, c=DARK)

tb(s, 0.55, 4.65, 4.30, 0.30, "AI 找 到 了 共同 点!",
   sz=11, b=True, c=PINK, a=PP_ALIGN.CENTER)

# RIGHT: explanation
panel(s, 5.15, 0.95, 4.45, 4.10, WARM, fill=WARM, lw=2)
tb(s, 5.30, 1.10, 4.15, 0.45, "AI 找: 「什么 一 样?」",
   sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.30, 1.58, 4.15, 0.32, "AI looks for what's similar",
   sz=11, c=GRAY, a=PP_ALIGN.LEFT)

# Explanation
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(2.10), Inches(4.15), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
ex.line.color.rgb = PINK; ex.line.width = Pt(2)
tb(s, 5.40, 2.20, 3.95, 0.30, "🎯 重点:", sz=11, b=True, c=PINK)
tb(s, 5.40, 2.50, 3.95, 0.35, "所有 猫 都 有 「尖 耳朵」 + 「胡须」!",
   sz=12, b=True, c=DARK)
tb(s, 5.40, 2.88, 3.95, 0.28, "All cats have these features",
   sz=9, c=GRAY)

# Kid analogy
ka = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(3.40), Inches(4.15), Inches(1.50))
ka.fill.solid(); ka.fill.fore_color.rgb = DAY
ka.line.color.rgb = STAR; ka.line.width = Pt(2.5)
tb(s, 5.40, 3.50, 3.95, 0.32, "💡 就 像 ...",
   sz=12, b=True, c=STAR)
tb(s, 5.40, 3.85, 3.95, 0.55, "你 知 道 苹果 都 是 圆 红 的!",
   sz=15, b=True, c=WHITE)
tb(s, 5.40, 4.45, 3.95, 0.32, "You know apples are round and red!",
   sz=10, c=WARM)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 关键 词: 「特 征 features」 — 帮 AI 分 类 的 共同 点\n• 让 学生 想: 你 怎么 一 眼 看 出 妈妈 是 妈妈? — 也 是 看 特征!\n• 高 年级 : AI 自 己 找 特征 — 不 是 人 教 的, 是 它 看 多 了 找 到 的")


# ============================================================
# 11 · STEP 5 — 测 试 (Test the AI with NEW data)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🧪 Step 5 · 测 试 · Test the AI", DAY)

# LEFT: test setup visual
panel(s, 0.40, 0.95, 4.55, 4.10, DAY, fill=WHITE, lw=3)
tb(s, 0.55, 1.05, 4.30, 0.50, "5️⃣", sz=28, b=True, c=DAY, a=PP_ALIGN.LEFT)

# New test image (NOT in training set)
ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.70), Inches(1.75), Inches(1.60), Inches(1.40))
ph.fill.solid(); ph.fill.fore_color.rgb = IMGBG
ph.line.color.rgb = DAY; ph.line.width = Pt(2)
tb(s, 0.70, 2.10, 1.60, 0.75, "🐱", sz=52, a=PP_ALIGN.CENTER)
tb(s, 0.70, 3.20, 1.60, 0.30, "测 试 图", sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)

# "Never seen" badge + arrow
tb(s, 2.45, 2.10, 0.55, 0.55, "❓", sz=30, a=PP_ALIGN.CENTER)
arrow(s, 2.45, 2.75, w=0.55, h=0.40, color=DAY)

# AI thinking
brain = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.15), Inches(1.75), Inches(1.60), Inches(1.40))
brain.fill.solid(); brain.fill.fore_color.rgb = WARM
brain.line.color.rgb = DAY; brain.line.width = Pt(3)
tb(s, 3.15, 1.95, 1.60, 0.85, "🤖💭", sz=46, a=PP_ALIGN.CENTER)
tb(s, 3.15, 3.20, 1.60, 0.30, "测 试 中 ...", sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)

tb(s, 0.55, 3.85, 4.30, 0.40, "用 「训练 时 没 见 过」 的 图 测!",
   sz=13, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.30, 4.30, 0.30, "训练 数据 ≠ 测 试 数据",
   sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.65, 4.30, 0.28, "Training data ≠ test data",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# RIGHT: explanation
panel(s, 5.15, 0.95, 4.45, 4.10, WARM, fill=WARM, lw=2)
tb(s, 5.30, 1.10, 4.15, 0.45, "测 试 AI 学 得 好 不 好!",
   sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.30, 1.58, 4.15, 0.32, "Test if AI learned well",
   sz=11, c=GRAY, a=PP_ALIGN.LEFT)

# Why test? — explanation
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(2.10), Inches(4.15), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
ex.line.color.rgb = DAY; ex.line.width = Pt(2)
tb(s, 5.40, 2.20, 3.95, 0.30, "🎯 为什么 测 试?", sz=11, b=True, c=DAY)
tb(s, 5.40, 2.50, 3.95, 0.35, "看 AI 能 不 能 用 在 「新」 情 况!",
   sz=12, b=True, c=DARK)
tb(s, 5.40, 2.88, 3.95, 0.28, "Check if AI generalizes to new cases",
   sz=9, c=GRAY)

# Kid analogy
ka = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(3.40), Inches(4.15), Inches(1.50))
ka.fill.solid(); ka.fill.fore_color.rgb = DAY
ka.line.color.rgb = STAR; ka.line.width = Pt(2.5)
tb(s, 5.40, 3.50, 3.95, 0.32, "💡 就 像 ...",
   sz=12, b=True, c=STAR)
tb(s, 5.40, 3.85, 3.95, 0.55, "老师 用 「新 题」 考 你!",
   sz=15, b=True, c=WHITE)
tb(s, 5.40, 4.45, 3.95, 0.32, "Like a teacher's test — with new questions!",
   sz=10, c=WARM)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 关键 概念: 「测 试 数据 vs 训练 数据」 — 必须 不同!\n• 用 训练 时 见 过 的 图 测试 = 作 弊! 不 能 真 检 验 AI\n• 真 正 严谨 的 测试: 用 全 新 的 图\n• 高 年级: 这 就 是 ML 的 train/test split 概念")


# ============================================================
# 12 · STEP 6 — 做 判 断 (AI Makes a Judgment with concrete examples)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "✨ Step 6 · 做 判 断 · Make a Judgment", GREEN)

# LEFT: 3 concrete examples of input → judgment
panel(s, 0.40, 0.95, 4.55, 4.10, GREEN, fill=WHITE, lw=3)
tb(s, 0.55, 1.05, 4.30, 0.40, "6️⃣  3 个 例子",
   sz=18, b=True, c=GREEN, a=PP_ALIGN.LEFT)
tb(s, 0.55, 1.45, 4.30, 0.25, "3 examples of AI judgments",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

# Three example rows: 输入 (input image) → 判断 (judgment)
examples = [
    ("🐱", "「这 是 猫!」", "✅", GREEN),
    ("🐶", "「这 是 狗!」", "✅", GREEN),
    ("🐰", "「不 知道...」", "❓", ORANGE),
]
for i, (input_em, judgment, mark, cl) in enumerate(examples):
    y = 1.85 + i * 1.05
    # Input box
    ib = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y), Inches(0.85), Inches(0.85))
    ib.fill.solid(); ib.fill.fore_color.rgb = IMGBG
    ib.line.color.rgb = cl; ib.line.width = Pt(1.5)
    tb(s, 0.55, y+0.05, 0.85, 0.75, input_em, sz=34, a=PP_ALIGN.CENTER)
    # Arrow
    arrow(s, 1.50, y+0.30, w=0.30, h=0.28, color=cl)
    # Judgment bubble
    jb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.90), Inches(y+0.08), Inches(2.35), Inches(0.70))
    jb.fill.solid(); jb.fill.fore_color.rgb = cl; jb.line.fill.background()
    tb(s, 1.95, y+0.10, 2.25, 0.25, "🤖", sz=11, a=PP_ALIGN.LEFT)
    tb(s, 1.95, y+0.32, 2.25, 0.45, judgment,
       sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Check / question mark
    tb(s, 4.35, y+0.20, 0.45, 0.50, mark, sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)

# Bottom note
tb(s, 0.55, 4.85, 4.30, 0.25, "💡 AI 见 过 的 → 答 对; 没 见 过 的 → 不 确定",
   sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

# RIGHT: explanation
panel(s, 5.15, 0.95, 4.45, 4.10, WARM, fill=WARM, lw=2)
tb(s, 5.30, 1.10, 4.15, 0.45, "AI 用 规律 做 判 断!",
   sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 5.30, 1.58, 4.15, 0.32, "AI uses learned patterns to decide",
   sz=11, c=GRAY, a=PP_ALIGN.LEFT)

# Explanation
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(2.10), Inches(4.15), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
ex.line.color.rgb = GREEN; ex.line.width = Pt(2)
tb(s, 5.40, 2.20, 3.95, 0.30, "🎯 完 成!", sz=11, b=True, c=GREEN)
tb(s, 5.40, 2.50, 3.95, 0.35, "学 得 越 好, 判 断 越 准!",
   sz=12, b=True, c=DARK)
tb(s, 5.40, 2.88, 3.95, 0.28, "Better learning = better judgments",
   sz=9, c=GRAY)

# Kid analogy
ka = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(3.40), Inches(4.15), Inches(1.50))
ka.fill.solid(); ka.fill.fore_color.rgb = DAY
ka.line.color.rgb = STAR; ka.line.width = Pt(2.5)
tb(s, 5.40, 3.50, 3.95, 0.32, "💡 就 像 ...",
   sz=12, b=True, c=STAR)
tb(s, 5.40, 3.85, 3.95, 0.55, "你 答 对 题 — 因为 你 学 好 了!",
   sz=15, b=True, c=WHITE)
tb(s, 5.40, 4.45, 3.95, 0.32, "Like answering correctly — because you studied well!",
   sz=10, c=WARM)
n += 1; pn(s, n)
notes(s, "3 分钟 — 用 具体 例子 解 释 「做 判 断」:\n• 例 1: AI 看 过 猫 → 给 它 猫 图 → 判 「猫!」 ✅\n• 例 2: AI 看 过 狗 → 给 它 狗 图 → 判 「狗!」 ✅\n• 例 3: AI 没 见 过 兔子 → 它 不 知道 怎么 判 → 可能 答 错 或 「不 确定」\n• 关 键: AI 只 能 判 它 「学 过」 的 东 西\n• 联 系: 这 也 是 为什么 我们 要 给 AI 多样 数据")


# ============================================================
# 13 · VIDEO — Machine Learning explainer (YouTube link)
# Teacher clicks the big PLAY button in slideshow mode to open the video.
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎬 看 视频 · Watch a Video!", DAY)

VIDEO_URL = "https://www.youtube.com/watch?v=_ZAz6_2TlcU"

tb(s, 0.4, 0.85, 9.2, 0.32, "学 完 6 步 — 现在 看 真 实 的 机器 学习!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "Now let's watch real ML in action!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: video thumbnail card (looks like a YouTube player)
thumb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.65), Inches(5.20), Inches(3.30))
thumb.fill.solid(); thumb.fill.fore_color.rgb = INK
thumb.line.color.rgb = RED; thumb.line.width = Pt(3)
# "YouTube" badge top-left
yb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.80), Inches(1.15), Inches(0.40))
yb.fill.solid(); yb.fill.fore_color.rgb = RED; yb.line.fill.background()
tb(s, 0.55, 1.83, 1.15, 0.32, "▶ YouTube", sz=11, b=True, c=WHITE, a=PP_ALIGN.CENTER)
# Big PLAY circle in center of thumbnail
play_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.30), Inches(2.50), Inches(1.40), Inches(1.40))
play_circle.fill.solid(); play_circle.fill.fore_color.rgb = RED
play_circle.line.color.rgb = WHITE; play_circle.line.width = Pt(4)
tb(s, 2.30, 2.75, 1.40, 0.90, "▶", sz=58, b=True, c=WHITE, a=PP_ALIGN.CENTER)
# Hyperlink on the play circle itself
play_circle.click_action.hyperlink.address = VIDEO_URL
# Caption below play button
tb(s, 0.55, 4.10, 4.90, 0.40, "🤖 Machine Learning 视频",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.50, 4.90, 0.30, "点 中 间 ▶ 按 钮 开 始 播 放!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)

# RIGHT: instructions + URL + reflection prompt
panel(s, 5.80, 1.65, 3.80, 3.30, DAY, fill=WARM, lw=3)

# Step-by-step play instructions
tb(s, 5.95, 1.78, 3.50, 0.35, "📺 怎么 看?",
   sz=14, b=True, c=DAY, a=PP_ALIGN.LEFT)
steps_play = [
    ("1.", "点 左 边 ▶ 按 钮"),
    ("2.", "视频 在 浏览 器 打 开"),
    ("3.", "全 班 一起 看"),
    ("4.", "看 完 一起 讨论!"),
]
for i, (num, txt) in enumerate(steps_play):
    y = 2.20 + i*0.38
    tb(s, 5.95, y, 0.30, 0.30, num, sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
    tb(s, 6.30, y, 3.20, 0.30, txt, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)

# URL display (so teachers can type if click fails)
url_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.95), Inches(3.80), Inches(3.55), Inches(0.55))
url_box.fill.solid(); url_box.fill.fore_color.rgb = WHITE
url_box.line.color.rgb = DAY; url_box.line.width = Pt(1.5)
tb(s, 6.05, 3.85, 3.40, 0.22, "🔗 网址:",
   sz=8, b=True, c=GRAY, a=PP_ALIGN.LEFT)
# URL as a hyperlinked run
url_tf_box = s.shapes.add_textbox(Inches(6.05), Inches(4.05), Inches(3.40), Inches(0.28))
url_tf = url_tf_box.text_frame
url_p = url_tf.paragraphs[0]
url_run = url_p.add_run()
url_run.text = "youtube.com/watch?v=_ZAz6_2TlcU"
url_run.font.size = Pt(9)
url_run.font.color.rgb = CYBER
url_run.font.name = 'KaiTi'
url_run.hyperlink.address = VIDEO_URL

# Reflection question
ref_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.95), Inches(4.42), Inches(3.55), Inches(0.50))
ref_box.fill.solid(); ref_box.fill.fore_color.rgb = STAR; ref_box.line.fill.background()
tb(s, 6.05, 4.47, 3.40, 0.40, "💬 看 完 想 一 想: AI 怎么 学 的?",
   sz=10, b=True, c=INK, a=PP_ALIGN.LEFT)

# Bottom takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.32))
tk.fill.solid(); tk.fill.fore_color.rgb = DAY; tk.line.fill.background()
tb(s, 0.55, 5.13, 9.0, 0.25, "🎯 看 视频 时 — 找 一 下 6 步!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5-8 分钟 视频 时 间:\n• 老师 在 演 示 时 点 ▶ 按 钮 — 视频 会 在 浏览 器 打 开\n• 视频 URL: https://www.youtube.com/watch?v=_ZAz6_2TlcU\n• 如果 链接 不 工作, 直接 在 浏览 器 输入 网 址\n• 看 视频 前 提示: 「看 一 下 你 能 不 能 找 出 我们 学 的 6 步」\n• 看 完 讨论: 视频 里 哪 些 是 我们 刚 学 的?\n• 备 案: 提前 下 载 视频 文件 以 防 网络 问题")



# ============================================================
# 10 · TRANSITION TO EXPERIMENT
# ============================================================
s = ns(prs); bg(s, INK, prs)
for x, y in [(0.4, 0.45), (9.1, 0.5), (0.5, 4.8), (9.0, 4.7), (1.2, 0.6), (8.3, 4.9)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.55, 9.4, 0.40, "🧪 实 验 时 间!",
   sz=18, b=True, c=NEON, a=PP_ALIGN.CENTER)

# Big question box
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.15), Inches(8.8), Inches(1.65))
qb.fill.solid(); qb.fill.fore_color.rgb = DAY
qb.line.color.rgb = STAR; qb.line.width = Pt(4)
tb(s, 0.8, 1.30, 8.4, 0.70, "AI 能 分 得 清 谁 是 谁 吗?",
   sz=30, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.00, 8.4, 0.40, "Can AI tell who's who?",
   sz=15, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.40, 8.4, 0.30, "今天 我们 来 训练 一 个 真 正 的 AI!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

# Lab equipment row
tb(s, 0.3, 3.20, 9.4, 1.10, "📷   🤖   🧠   🧪   ✨",
   sz=58, a=PP_ALIGN.CENTER)

# TM banner
bn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.50), Inches(8.8), Inches(0.95))
bn.fill.solid(); bn.fill.fore_color.rgb = WHITE
bn.line.color.rgb = STAR; bn.line.width = Pt(3)
tb(s, 0.75, 4.60, 8.5, 0.40, "💻 用 Teachable Machine — 一 起 训练 AI!",
   sz=15, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.75, 5.05, 8.5, 0.30, "Let's train an AI together — live!",
   sz=11, c=GRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "30 秒 hype slide.\n• 老师 切到 投影 — 显示 teachablemachine.withgoogle.com\n• 大 声 宣布: 「今天 我们 不 是 看 视频, 是 做 实验!」\n• 让 学生 鼓掌, 兴奋 起来")


# ============================================================
# 11 · CORE EXPERIMENT SETUP — 大咪 vs 小咪
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🐱 大 咪 vs 小 咪 · Cat A vs Cat B", DAY)

# Top description
tb(s, 0.4, 0.85, 9.2, 0.32, "我们 要 训练 AI 认 两 只 长 得 像 的 猫!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Two photo placeholders side by side
photo_slot(s, 0.50, 1.30, 4.30, 2.70, "📷 大 咪 照片",
           "Tabby + white paws", DAY)
photo_slot(s, 5.20, 1.30, 4.30, 2.70, "📷 小 咪 照片",
           "Tabby, NO white paws", ORANGE)

# Cat descriptions below photos
desc1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(4.05), Inches(4.30), Inches(0.55))
desc1.fill.solid(); desc1.fill.fore_color.rgb = DAY; desc1.line.fill.background()
tb(s, 0.55, 4.10, 4.20, 0.30, "🐱 大 咪 Da-mi",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.38, 4.20, 0.22, "虎斑 + 白色 脚 印 + 特别 花纹",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)

desc2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.20), Inches(4.05), Inches(4.30), Inches(0.55))
desc2.fill.solid(); desc2.fill.fore_color.rgb = ORANGE; desc2.line.fill.background()
tb(s, 5.25, 4.10, 4.20, 0.30, "🐱 小 咪 Xiao-mi",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 5.25, 4.38, 4.20, 0.22, "虎斑 + 没 有 白 脚 + 灰白 身 体",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)

# Big prediction question
q = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.65))
q.fill.solid(); q.fill.fore_color.rgb = INK
q.line.color.rgb = STAR; q.line.width = Pt(2)
tb(s, 0.55, 4.82, 9.0, 0.32, "🤔 AI 能 分 得 清 大 咪 和 小 咪 吗? 投 票!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.13, 9.0, 0.26, "Will AI tell them apart? Vote: 能 / 不 能 / 不 知 道",
   sz=9, c=LGRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 分钟 setup:\n• 老师 提前 拍 好 2 只 课堂 猫 (或 用 学生 自带 玩 偶) — 备 好 照片\n• 也 可以 用 老师 家 的 真 猫 照片\n• 重点: 两 只 猫 要 「长 得 像 但 有 差别」 (e.g., 都 是 虎 斑, 一 只 有 白 脚 一 只 没 有)\n• 让 学生 先 看 + 比较 + 找 特征\n• 然后 投票 预测 — 制造 悬念!")


# ============================================================
# 12 · EXPERIMENT ROUND 1 — Few photos
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔬 实验 第 1 轮 · Round 1: Few Photos", PINK)

# Setup card
setup = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.75))
setup.fill.solid(); setup.fill.fore_color.rgb = WARM
setup.line.color.rgb = PINK; setup.line.width = Pt(2.5)
tb(s, 0.55, 1.02, 9.0, 0.32, "📷 训练 数据: 每 只 猫 只 拍 1-2 张!",
   sz=15, b=True, c=PINK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.38, 9.0, 0.28, "Train data: only 1-2 photos per cat",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Visual: very few photos
panel(s, 0.40, 1.85, 4.40, 2.30, DAY, fill=WHITE, lw=2.5)
panel_head(s, 0.40, 1.85, 4.40, DAY, "🐱 大 咪 — 只 有 1 张", sz=12)
tb(s, 0.50, 2.45, 4.20, 1.60, "📷", sz=110, a=PP_ALIGN.CENTER)

panel(s, 5.20, 1.85, 4.40, 2.30, ORANGE, fill=WHITE, lw=2.5)
panel_head(s, 5.20, 1.85, 4.40, ORANGE, "🐱 小 咪 — 只 有 1 张", sz=12)
tb(s, 5.30, 2.45, 4.20, 1.60, "📷", sz=110, a=PP_ALIGN.CENTER)

# Live test prompt — shorter box so text fits within slide bounds
test = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.40), Inches(9.20), Inches(1.00))
test.fill.solid(); test.fill.fore_color.rgb = PINK
test.line.color.rgb = STAR; test.line.width = Pt(3)
tb(s, 0.55, 4.50, 9.0, 0.45, "🧪 现场 测试! AI 会 猜 对 吗?",
   sz=20, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.00, 9.0, 0.30, "Live test! 👉 一起 投 票: 能 / 不 能",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5-7 分钟 现场 演示 (核心 体验!):\n• 老师 在 Teachable Machine 上 只 拍 1-2 张 大 咪 + 1-2 张 小 咪\n• 点 Train → 立 即 测试\n• 给 它 看 大 咪 新 照片 → 看 AI 怎么 说\n• 多 半 会 错! 因为 数据 太 少\n• 让 学生 大 声 喊 出 AI 的 预测\n• 故意 制造 错误 — 这 是 下 一 张 的 讨论 引子")


# ============================================================
# 13 · DISCUSSION — Why AI got it wrong
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🤔 为什么 AI 猜 错 了? · Why Did AI Get It Wrong?", PINK)

tb(s, 0.4, 0.85, 9.2, 0.32, "想 一 想 — AI 为什么 分 不 清?",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 reason cards in 2x2
reasons = [
    ("📷", "图片 太 少", "Too few images", CYBER),
    ("📚", "学 得 不 够", "Didn't learn enough", ORANGE),
    ("🔍", "不 知 道 特征", "Doesn't know features", PURPLE),
    ("🐱", "长 得 太 像", "Cats look too similar", PINK),
]
card_w = 4.30; card_h = 1.65; gap_x = 0.25; gap_y = 0.20
start_x = (10 - 2*card_w - gap_x)/2
for i, (em, cn, en, cl) in enumerate(reasons):
    row = i // 2; col = i % 2
    x = start_x + col*(card_w + gap_x)
    y = 1.30 + row*(card_h + gap_y)
    panel(s, x, y, card_w, card_h, cl, fill=WHITE, lw=3)
    tb(s, x+0.20, y+0.30, 0.90, 0.90, em, sz=42, a=PP_ALIGN.LEFT)
    tb(s, x+1.25, y+0.35, card_w-1.40, 0.50, cn, sz=18, b=True, c=cl)
    tb(s, x+1.25, y+0.90, card_w-1.40, 0.40, en, sz=10, c=GRAY)

# Bottom prompt
tb(s, 0.4, 5.10, 9.2, 0.30, "💬 小 组 讨 论 — 还 有 其他 原因 吗?",
   sz=11, b=True, c=PINK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3-5 分钟 讨论:\n• 让 学生 大 声 说 出 想 法\n• 引导 — 把 答案 落 到 「数据 少」 + 「特征 难 抓」\n• 这 是 关键 教学 时刻: AI 不 是 不 聪明, 是 「学 习 材料」 不 够")


# ============================================================
# 14 · EXPERIMENT ROUND 2 — More photos
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔬 实验 第 2 轮 · Round 2: More Photos!", GREEN)

# Setup card
setup = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.75))
setup.fill.solid(); setup.fill.fore_color.rgb = WARM
setup.line.color.rgb = GREEN; setup.line.width = Pt(2.5)
tb(s, 0.55, 1.02, 9.0, 0.32, "📷 加 多 数据! 每 只 猫 5-10 张 — 不同 角度!",
   sz=15, b=True, c=GREEN, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.38, 9.0, 0.28, "Add more data! 5-10 photos per cat — different angles!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Variety tags
varieties = [
    ("📐", "不同 角度"),
    ("🤸", "不同 姿势"),
    ("💡", "不同 光线"),
    ("📏", "不同 距离"),
]
vw = 2.10; vgap = 0.15
vtotal = 4*vw + 3*vgap; vstart = (10 - vtotal)/2
for i, (em, txt) in enumerate(varieties):
    x = vstart + i*(vw + vgap)
    panel(s, x, 2.00, vw, 1.05, GREEN, fill=WHITE, lw=2.5)
    tb(s, x, 2.10, vw, 0.50, em, sz=26, a=PP_ALIGN.CENTER)
    tb(s, x, 2.62, vw, 0.40, txt, sz=12, b=True, c=GREEN, a=PP_ALIGN.CENTER)

# Many photos visual
tb(s, 0.4, 3.20, 9.2, 0.85, "📷 📷 📷 📷 📷  vs  📷 📷 📷 📷 📷",
   sz=36, a=PP_ALIGN.CENTER)
tb(s, 0.4, 4.05, 9.2, 0.30, "10 张 大 咪 + 10 张 小 咪",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Live retest prompt
retest = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.50), Inches(9.20), Inches(0.90))
retest.fill.solid(); retest.fill.fore_color.rgb = GREEN
retest.line.color.rgb = STAR; retest.line.width = Pt(3)
tb(s, 0.55, 4.62, 9.0, 0.40, "🧪 再 测 试! 现在 AI 会 更 聪 明 吗?",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.05, 9.0, 0.30, "Retest now! Will AI be smarter?",
   sz=11, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5-7 分钟 现场 演示:\n• 老师 加 多 照片 — 拍 不同 角度 的 大 咪 + 小 咪\n• 让 学生 帮 忙 决定 拍 什么 (头 / 尾巴 / 全 身)\n• 重新 训练 → 测试\n• 这 次 应 该 答 对 大 部分\n• 让 学生 鼓掌!\n• 引出 下 一 张: 为什么 这 次 好 了?")


# ============================================================
# 15 · DISCUSSION — Why it got better
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "✨ 为什么 现在 AI 更 棒? · Why Did AI Improve?", GREEN)

# 3 reasons
reasons = [
    ("📚", "更 多 数据", "More data", CYBER),
    ("👀", "看 过 更 多 例子", "Saw more examples", PINK),
    ("🔍", "学 到 更 多 特征", "Learned more features", ORANGE),
]
card_w = 2.85; gap = 0.20
total = 3*card_w + 2*gap; start = (10 - total)/2
for i, (em, cn, en, cl) in enumerate(reasons):
    x = start + i*(card_w + gap)
    panel(s, x, 1.05, card_w, 2.85, cl, fill=WHITE, lw=3)
    tb(s, x, 1.20, card_w, 0.95, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 2.25, card_w, 0.50, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 2.78, card_w, 0.32, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.25, card_w-0.20, 0.55, str(i+1), sz=40, b=True, c=cl, a=PP_ALIGN.CENTER)

# Big takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.20), Inches(9.20), Inches(1.20))
tk.fill.solid(); tk.fill.fore_color.rgb = GREEN
tk.line.color.rgb = STAR; tk.line.width = Pt(4)
tb(s, 0.55, 4.35, 9.0, 0.55, "💡 更 多 数据 = 更 好 学习!",
   sz=28, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.95, 9.0, 0.35, "More data = better learning!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 分钟:\n• 让 学生 一起 喊 「更 多 数据 = 更 好 学习!」\n• 这 是 今天 最 重 要 的 一 句 话\n• 跟 第 5 张 (baby) 呼应: 看 多 了 就 学 会")


# ============================================================
# 16 · CHALLENGE ROUND — Tricky tests
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎯 大 挑 战! · Challenge Round!", ORANGE)

tb(s, 0.4, 0.85, 9.2, 0.32, "用 「难」 的 图 测试 AI — 它 还 认 得 吗?",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "Test AI with tricky images — can it still tell?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 4 tricky challenges
tricky = [
    ("📐", "角度 奇 怪", "Unusual angle", "倒 着 拍?"),
    ("🌫️", "模糊 的 图", "Blurry photo", "看 不 清?"),
    ("🐈", "长 得 像 的 别 的 猫", "Similar other cat", "不 是 大 咪 / 小 咪!"),
    ("🧸", "玩 具 猫", "Toy cat", "是 真 的 吗?"),
]
card_w = 2.15; gap = 0.20
total = 4*card_w + 3*gap; start = (10 - total)/2
for i, (em, cn, en, ex) in enumerate(tricky):
    x = start + i*(card_w + gap)
    panel(s, x, 1.65, card_w, 2.85, ORANGE, fill=WHITE, lw=3)
    tb(s, x, 1.78, card_w, 0.85, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, card_w, 0.40, cn, sz=14, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
    tb(s, x, 3.12, card_w, 0.28, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.08, 3.55, card_w-0.16, 0.80, ex, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Big question
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.70), Inches(9.20), Inches(0.75))
qb.fill.solid(); qb.fill.fore_color.rgb = DAY
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 4.78, 9.0, 0.40, "🤔 AI 还 认 得 吗?",
   sz=20, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.20, 9.0, 0.28, "Can AI still recognize?",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 演示:\n• 用 这 4 类 「难」 图 测试 已 经 训练 好 的 AI\n• 多 半 又 会 错!\n• 让 学生 看 到: AI 不 是 万 能, 它 只 学 了 「正常」 的 照片\n• 引出 下 一 张: AI 也 会 犯错")


# ============================================================
# 17 · AI MAKES MISTAKES
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "😅 AI 也 会 犯 错! · AI Makes Mistakes!", PINK)

tb(s, 0.4, 0.85, 9.2, 0.32, "AI 可能 会 犯错, 因为:",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 3 mistake reasons
mistakes = [
    ("📉", "数据 太 少", "Too little data", CYBER),
    ("🔍", "特征 太 像", "Features too similar", ORANGE),
    ("🌫️", "图片 不 清 楚", "Image unclear", PURPLE),
]
card_w = 2.65; gap = 0.30
total = 3*card_w + 2*gap; start = (10 - total)/2
for i, (em, cn, en, cl) in enumerate(mistakes):
    x = start + i*(card_w + gap)
    panel(s, x, 1.25, card_w, 2.85, cl, fill=WHITE, lw=3)
    tb(s, x, 1.38, card_w, 0.85, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 2.30, card_w, 0.45, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 2.78, card_w, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+card_w/2-0.30), Inches(3.35), Inches(0.60), Inches(0.60))
    badge.fill.solid(); badge.fill.fore_color.rgb = cl; badge.line.fill.background()
    tb(s, x+card_w/2-0.30, 3.40, 0.60, 0.50, str(i+1), sz=22, b=True, c=WHITE, a=PP_ALIGN.CENTER)

# Big message
mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.35), Inches(9.20), Inches(1.05))
mb.fill.solid(); mb.fill.fore_color.rgb = PINK
mb.line.color.rgb = STAR; mb.line.width = Pt(3)
tb(s, 0.55, 4.50, 9.0, 0.50, "🤖 AI 不 是 完美 的 — 它 在 学习!",
   sz=22, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.05, 9.0, 0.30, "AI is not perfect — it's still learning!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 分钟:\n• 让 学生 明白: AI 也 会 错 是 「正常」 的\n• 关键: 我们 怎么 帮 AI 变 更 聪明? → 给 它 更 多 + 更 好 数据!\n• 高 年级: 可以 引入 「测试 数据 vs 训练 数据」 概念")


# ============================================================
# SMALL GROUP REFLECTION
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "💬 小 组 讨 论 · Group Discussion", DAY)

# Big question
q = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.15))
q.fill.solid(); q.fill.fore_color.rgb = DAY
q.line.color.rgb = STAR; q.line.width = Pt(3)
tb(s, 0.55, 1.10, 9.0, 0.50, "🤔 如果 你 要 训练 AI, 你 会 给 它 什么 样 的 图片?",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.65, 9.0, 0.32, "If YOU train an AI — what photos would you give it?",
   sz=11, c=WARM, a=PP_ALIGN.CENTER)

# 5 idea bubbles in row
ideas = [
    ("📐", "不同 角度", "Different angles"),
    ("📷", "清楚 的 图", "Clear photos"),
    ("📏", "远 + 近", "Near + far"),
    ("🌈", "不同 背景", "Different backgrounds"),
    ("💯", "很 多 例子", "Lots of examples"),
]
iw = 1.75; igap = 0.10
itotal = 5*iw + 4*igap; istart = (10 - itotal)/2
for i, (em, cn, en) in enumerate(ideas):
    x = istart + i*(iw + igap)
    panel(s, x, 2.40, iw, 2.20, ORANGE, fill=WHITE, lw=2.5)
    tb(s, x, 2.55, iw, 0.75, em, sz=38, a=PP_ALIGN.CENTER)
    tb(s, x, 3.40, iw, 0.50, cn, sz=12, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.95, iw-0.10, 0.35, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# Group discussion prompt
disc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.80), Inches(9.20), Inches(0.60))
disc.fill.solid(); disc.fill.fore_color.rgb = INK; disc.line.fill.background()
tb(s, 0.55, 4.88, 9.0, 0.32, "👥 4 人 一 组 — 3 分钟 讨论 后 分 享!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.18, 9.0, 0.22, "Groups of 4 — discuss 3 min, then share!",
   sz=9, c=LGRAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟:\n• 小 组 讨论 — 让 学生 自己 想 答案 (不 提示)\n• 然后 分 享 — 老师 收 集 答 案 写 白 板 上\n• 这 是 「学 当 一 个 AI 训练 师」 的 关键 时刻\n• 让 学生 感受 自己 是 设 计 者 不 是 用 户")


# ============================================================
# 20 · WRAP-UP — repeat the same ML 6 steps (consistent terminology)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎓 总 结 · Wrap-up · 机器 学习 6 步", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "今天 我们 学 了 机器 学习 的 6 个 步 骤 — 一起 复 习!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Same 6 steps as overview (slide 6) — KEEP CONSISTENT
wrap_steps = [
    ("1️⃣", "📷", "收集 例子", "Collect", CYBER),
    ("2️⃣", "🏷️", "贴 标签", "Label", ORANGE),
    ("3️⃣", "🧠", "AI 学 习", "Study", PURPLE),
    ("4️⃣", "🔍", "找 规律", "Patterns", PINK),
    ("5️⃣", "🧪", "测 试", "Test", DAY),
    ("6️⃣", "✨", "做 判 断", "Predict", GREEN),
]
ww = 2.85; wgap_x = 0.20; wgap_y = 0.20; wh = 1.55
wstart_x = (10 - 3*ww - 2*wgap_x)/2
for i, (num, em, cn, en, cl) in enumerate(wrap_steps):
    row = i // 3; col = i % 3
    x = wstart_x + col*(ww + wgap_x)
    y = 1.30 + row*(wh + wgap_y)
    panel(s, x, y, ww, wh, cl, fill=WHITE, lw=2.5)
    tb(s, x+0.10, y+0.08, 0.55, 0.45, num, sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x+0.60, y+0.08, ww-0.70, 0.70, em, sz=38, a=PP_ALIGN.CENTER)
    tb(s, x, y+0.85, ww, 0.38, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, y+1.20, ww, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Final takeaway bar
fin = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.50))
fin.fill.solid(); fin.fill.fore_color.rgb = DAY
fin.line.color.rgb = STAR; fin.line.width = Pt(3)
tb(s, 0.55, 5.05, 9.0, 0.35, "💡 机器 学习 不 是 魔法 — 是 6 步 学 出 来 的!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 分钟 总 结:\n• 全 班 一 起 数 「1, 2, 3, 4, 5, 6 步!」\n• 一 起 念 每 一 步 的 中文 名 字\n• 跟 Slide 6 (开 头 overview) 一 样 — 帮 学 生 加 深 印象\n• 全 班 一 起 大 声 喊: 「机器 学习 不 是 魔法 — 是 6 步 学 出 来 的!」\n• 下 一: Session 2 学 词汇 / Session 3 自己 训练 AI")


# ============================================================
# 9 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2", "📖 下午 2:00–2:45  ·  复习 + 中文 词汇 · 我 会 认 / 我 会 写",
        DAY, "📚"); n += 1; pn(s, n)


# ============================================================
# 10-13 · 我 会 认 (vocabulary recognition)
# ============================================================
recognize_words = [
    ("🤖", "机器", "jī qì", "Machine",
     "机器 人 帮 我们 做 事。", "Robots help us do things.",
     "机器 人 / 机械 手 / 工厂 机器", CYBER),
    ("🗂️", "分类", "fēn lèi", "Classify",
     "我们 给 图片 分类。", "We sort the pictures.",
     "整理 文件夹 / 分 类 物品", ORANGE),
    ("📊", "数据", "shù jù", "Data",
     "电脑 需要 很 多 数据 学 习。", "Computers need lots of data to learn.",
     "图表 / 数字 / 表 格", PINK),
    ("🏋️", "训练", "xùn liàn", "Train",
     "我们 要 训练 一 个 AI。", "We're training an AI.",
     "训练 中 / 运动 训练 / 教 AI", DAY),
]
for em, cn, py, en, ex_cn, ex_en, hint, cl in recognize_words:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)


# ============================================================
# 14-15 · 我 会 写 (writing practice)
# ============================================================
s = vocab_write(prs, ORANGE, "分类", "Classify",
                [("分", "fēn", "4 笔", "上 「八」 下 「刀」 — 切 开 分 开"),
                 ("类", "lèi", "9 笔", "上 「米」 下 「大」")])
n += 1; pn(s, n)

s = vocab_write(prs, DAY, "训练", "Train",
                [("训", "xùn", "5 笔", "「讠」 + 「川」 — 用 话 教"),
                 ("练", "liàn", "8 笔", "「纟」 + 「东」 — 反复 练 习")])
n += 1; pn(s, n)


# ============================================================
# 16 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3", "🎨 下午 3:00–4:30  ·  训练 真 AI + 我 是 AI 训练师!",
        DAY, "🤖"); n += 1; pn(s, n)

# Complete today's booklet — before project starts
s = booklet_slide(prs, day_num=3, day_topic_cn="机器 学习 · 我 是 AI 训练 师", day_color=DAY)
n += 1; pn(s, n)


# ============================================================
# TM 5-STEP INSTRUCTIONS + THEMES  (Session 3 opener — hands-on prep)
# Note: Cat-differentiation example was MOVED to Session 1 (Rounds 1 & 2).
# Session 3 now starts directly with the hands-on instructions for student groups.
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📝 操作 步骤 · How to Use Teachable Machine", DAY)

# LEFT: 5 steps
panel(s, 0.40, 0.95, 4.55, 3.95, DAY, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, DAY, "5 步 操作  5 Steps", sz=14)
steps = [
    ("1️⃣", "选择 分类 主题", "Pick a theme"),
    ("2️⃣", "收集 图片", "Collect images"),
    ("3️⃣", "点击 训练", "Click Train"),
    ("4️⃣", "测试 AI", "Test the AI"),
    ("5️⃣", "看 AI 会 不 会 错!", "See if AI gets it wrong!"),
]
for i, (num, cn, en) in enumerate(steps):
    y = 1.60 + i * 0.62
    tb(s, 0.55, y, 0.55, 0.50, num, sz=20, b=True, c=DAY, a=PP_ALIGN.LEFT)
    tb(s, 1.15, y+0.02, 3.70, 0.32, cn, sz=13, b=True, c=DARK)
    tb(s, 1.15, y+0.33, 3.70, 0.25, en, sz=9, c=GRAY)

# RIGHT: 5 themes
panel(s, 5.05, 0.95, 4.55, 3.95, ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.05, 0.95, 4.55, ORANGE, "🎨 主题 选择  Themes", sz=14)
themes = [
    ("🐾", "动 物", "Animals"),
    ("🍎", "水 果", "Fruits"),
    ("😊", "表 情", "Emotions"),
    ("✋", "手 势", "Hand gestures"),
    ("🌈", "颜色 物品", "Colored objects"),
]
for i, (em, cn, en) in enumerate(themes):
    y = 1.60 + i * 0.62
    tb(s, 5.20, y, 0.55, 0.50, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, 5.85, y+0.02, 3.55, 0.32, cn, sz=14, b=True, c=ORANGE)
    tb(s, 5.85, y+0.33, 3.55, 0.25, en, sz=9, c=GRAY)

url = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.42))
url.fill.solid(); url.fill.fore_color.rgb = INK; url.line.fill.background()
tb(s, 0.55, 5.05, 9.0, 0.32, "💻 teachablemachine.withgoogle.com  ·  免费 + 不 用 注册!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "操作 建议:\n• 全 班 一起 看 投影 演示 一 次\n• 然后 小 组 用 iPad 自己 玩\n• K-2: 选 简单 的 (动物 / 水果)\n• 3-5: 可以 尝试 手势 / 表情\n• 重要: 让 学生 自己 拍 照, 不 要 老师 代 劳")


# ============================================================
# 19 · BAD DATA EXAMPLE (错 数据 = 错 结果)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "😅 错 数据 = 错 结果! · Bad Data = Bad Results", PINK)

# LEFT: training panel (all white shapes labeled "rabbit")
panel(s, 0.40, 0.95, 4.30, 3.80, ORANGE, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.30, ORANGE, "📚 训练: 全 部 是 白色!", sz=13)

train = [
    ("白 兔子", "→ 「兔子」 ✅"),
    ("白 狗 狗", "→ 「兔子」 ⁉️"),
    ("白 猫 咪", "→ 「兔子」 ⁉️"),
]
for i, (label, pred) in enumerate(train):
    y = 1.65 + i * 1.00
    oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.60), Inches(y), Inches(0.85), Inches(0.85))
    oval.fill.solid(); oval.fill.fore_color.rgb = WHITE
    oval.line.color.rgb = DARK; oval.line.width = Pt(2)
    tb(s, 0.60, y+0.22, 0.85, 0.45, "白", sz=22, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 1.55, y+0.10, 1.20, 0.35, label, sz=13, b=True, c=DARK)
    arrow(s, 2.78, y+0.32, w=0.28, h=0.25, color=ORANGE)
    tb(s, 3.10, y+0.20, 1.60, 0.50, pred, sz=11, b=True, c=ORANGE, a=PP_ALIGN.LEFT)

arrow(s, 4.80, 2.60, w=0.45, h=0.45, color=PINK)

# RIGHT: test panel
panel(s, 5.30, 0.95, 4.30, 3.80, PINK, fill=WHITE, lw=3)
panel_head(s, 5.30, 0.95, 4.30, PINK, "🧪 测试: 给 AI 看 新 图", sz=13)
big = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.45), Inches(1.65), Inches(2.00), Inches(1.55))
big.fill.solid(); big.fill.fore_color.rgb = WHITE
big.line.color.rgb = DARK; big.line.width = Pt(3)
tb(s, 6.45, 2.05, 2.00, 0.60, "白", sz=36, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 5.30, 3.30, 4.30, 0.32, "白色 北极熊  Polar Bear",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
ai = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(3.78), Inches(3.80), Inches(0.90))
ai.fill.solid(); ai.fill.fore_color.rgb = PINK; ai.line.fill.background()
tb(s, 5.65, 3.85, 3.60, 0.28, "🤖 AI 说:",
   sz=11, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 5.65, 4.15, 3.60, 0.50, "「兔子!」",
   sz=26, b=True, c=WHITE, a=PP_ALIGN.CENTER)

tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.55))
tk.fill.solid(); tk.fill.fore_color.rgb = INK; tk.line.fill.background()
tb(s, 0.55, 4.92, 9.0, 0.40, "💡 错 数据 = 错 结果!  AI 只 学 到 「白色 = 兔子」!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "搞笑 + 重要!\n• 让 学生 一起 喊 「兔子!」 — 加强 记忆\n• 训练 数据 全部 是 白色 → AI 学 到 「白色 = 兔子」\n• 给 它 新 的 白色 东西 (北极熊) → 它 还是 说 「兔子」\n• 真实 例子: 如果 训练 数据 不 多 样, AI 就 不 公平\n• 这 是 AI 公平性 (data bias) 的 入门 话题")


# ============================================================
# 20 · THE MORAL — AI 像 你 in school
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🏫 AI 学 习 像 你 在 学校! · AI Learns Like You at School", DAY)

# Big banner
bn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.75))
bn.fill.solid(); bn.fill.fore_color.rgb = WARM
bn.line.color.rgb = DAY; bn.line.width = Pt(2.5)
tb(s, 0.55, 1.02, 9.0, 0.35, "好 好 学 = 好 结果! 你 这 样, AI 也 一 样!",
   sz=15, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.38, 9.0, 0.28, "Study well = good results. Same for you AND for AI.",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 2 columns: Student | AI
cols = [
    ("🎒", "你 (学 生) · You", DAY, [
        ("✅", "认真 听 课 + 做 作业", "Good: pay attention + do homework"),
        ("🏆", "→ 考 试 答 对!", "→ Get answers right!"),
        ("❌", "走 神 + 不 复习", "Bad: zone out + don't review"),
        ("😅", "→ 考 试 答 错!", "→ Get answers wrong!"),
    ]),
    ("🤖", "AI · AI", PINK, [
        ("✅", "好 数据 + 多 例子", "Good: clean data + lots of examples"),
        ("🏆", "→ AI 答 对!", "→ AI predicts correctly!"),
        ("❌", "坏 数据 + 错 标签", "Bad: wrong data + bad labels"),
        ("😅", "→ AI 答 错!", "→ AI gets it wrong!"),
    ]),
]
col_w = 4.40; gap = 0.40
total = 2*col_w + gap; start = (10 - total)/2
for i, (em, title, cl, bullets) in enumerate(cols):
    x = start + i*(col_w + gap)
    panel(s, x, 1.85, col_w, 2.95, cl, fill=WHITE, lw=3)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85), Inches(col_w), Inches(0.50))
    head.fill.solid(); head.fill.fore_color.rgb = cl; head.line.fill.background()
    tb(s, x+0.10, 1.92, col_w-0.20, 0.40, f"{em}  {title}",
       sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    for j, (e, cn, en) in enumerate(bullets):
        y = 2.45 + j*0.58
        tb(s, x+0.18, y, 0.40, 0.45, e, sz=18, a=PP_ALIGN.LEFT)
        tb(s, x+0.65, y+0.04, col_w-0.80, 0.28, cn, sz=12, b=True, c=DARK)
        tb(s, x+0.65, y+0.32, col_w-0.80, 0.24, en, sz=8, c=GRAY)

# Bottom moral
mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.45))
mb.fill.solid(); mb.fill.fore_color.rgb = DAY
mb.line.color.rgb = STAR; mb.line.width = Pt(2)
tb(s, 0.55, 5.00, 9.0, 0.35, "💡 给 自己 好 的 学 习, 给 AI 好 的 数据 — 我们 都 是 学 习 者!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "核心 教学 时刻!\n• 让 学生 自己 说: AI 学 习 像 我 在 学校 — 怎么 像?\n• 引导 出: 都 需要 「好 例子」 + 「多 练 习」\n• 数据 公平 = 数据 多 样 (不 只 是 一 种)\n• 这 是 道德 + STEM 结 合 — 让 学生 思考: 怎么 当 「好 AI 训练师」?\n• 联 想: 我 们 平 时 给 自己 看 什么? 玩 游戏? 还是 学 知识?")


# ============================================================
# EMOTION TRAINING DATA — same logic as morning cat experiment
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎭 用 表情 照片 训练 AI · Train AI on Emotions", DAY)

# Bridge: cat → emotion (same 6 steps!)
bridge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.70))
bridge.fill.solid(); bridge.fill.fore_color.rgb = WARM
bridge.line.color.rgb = DAY; bridge.line.width = Pt(2.5)
tb(s, 0.55, 1.02, 9.0, 0.32, "🐱 上午: 训练 AI 认 猫  →  🎭 现在: 训练 AI 认 表情",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.34, 9.0, 0.28, "Same 6 steps — different topic!  和 上午 一 样 的 6 步!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 4 emotion training data columns — each shows multiple training photos
emotions = [
    ("开 心", "Happy", "😊  😄\n😁  🙂", GREEN),
    ("生 气", "Angry", "😡  😠\n🤬  😤", RED),
    ("难 过", "Sad", "😢  😭\n😞  😔", CYBER),
    ("惊 讶", "Surprised", "😲  😮\n😱  🤯", PURPLE),
]
ew = 2.15; egap = 0.18
etotal = 4*ew + 3*egap; estart = (10 - etotal)/2
for i, (cn, en, photos, cl) in enumerate(emotions):
    x = estart + i*(ew + egap)
    panel(s, x, 1.85, ew, 2.85, cl, fill=WHITE, lw=2.5)
    # Header: emotion label
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85), Inches(ew), Inches(0.55))
    head.fill.solid(); head.fill.fore_color.rgb = cl; head.line.fill.background()
    tb(s, x, 1.93, ew, 0.40, cn, sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # "训练 数据" label
    tb(s, x, 2.50, ew, 0.25, "📷 训练 数据:", sz=9, b=True, c=cl, a=PP_ALIGN.CENTER)
    # Multiple emotion photos (training examples)
    bx = s.shapes.add_textbox(Inches(x), Inches(2.80), Inches(ew), Inches(1.55))
    tf = bx.text_frame; tf.word_wrap = True
    lines = photos.split("\n")
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = lines[0]; r0.font.size = Pt(36); r0.font.name = 'KaiTi'
    for line in lines[1:]:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line; r.font.size = Pt(36); r.font.name = 'KaiTi'
    # English subtitle
    tb(s, x, 4.42, ew, 0.25, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom callout: same 6 steps
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.60))
cb.fill.solid(); cb.fill.fore_color.rgb = DAY
cb.line.color.rgb = STAR; cb.line.width = Pt(2.5)
tb(s, 0.55, 4.93, 9.0, 0.30, "💡 同 样 的 6 步: 收集 → 标签 → 学习 → 找规律 → 测试 → 做判断",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22, "Same 6 steps as the morning cat experiment!",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3-5 分钟 (桥 接 上午 + 下午 项目):\n• 重 点: 让 学 生 看 到 「同 一 个 ML 流程 可以 学 不同 东 西」\n• 上午: 给 AI 看 100 张 猫 照片 → AI 认 猫\n• 现在: 给 AI 看 100 张 开心 照片 → AI 认 开心\n• 类 别 不 同, 但 6 步 完 全 一 样\n• 老师 可以 用 投 影 演示: 在 Teachable Machine 上 加 几 张 自己 的 表情 照片 当 演示\n• 引出 下 一 张: 「现在 我们 自己 设 计 一 个 表情 识别 机器人!」")


# ============================================================
# PROJECT SLIDE 1 · TRANSITION — recap Session 1 → new robot project
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎯 刚才 我们 训练 了 AI! · We Just Trained AI!", DAY)

# LEFT: recap from Session 1
panel(s, 0.40, 0.95, 4.40, 3.95, CYBER, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.40, CYBER, "🐱 早上 你 做 了 什么?", sz=12)
tb(s, 0.55, 1.65, 4.10, 0.85, "🐱  🔍",
   sz=58, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.55, 4.10, 0.40, "训练 AI 认 「猫」",
   sz=16, b=True, c=CYBER, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.95, 4.10, 0.30, "Trained AI to recognize cats",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)
# Features recap
fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(3.40), Inches(4.10), Inches(1.40))
fr.fill.solid(); fr.fill.fore_color.rgb = WARM; fr.line.fill.background()
tb(s, 0.65, 3.48, 3.90, 0.30, "🎯 AI 看 什么 特 征?",
   sz=12, b=True, c=CYBER, a=PP_ALIGN.LEFT)
tb(s, 0.65, 3.80, 3.90, 0.32, "✓ 尖 耳朵   ✓ 胡须   ✓ 尾巴",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.65, 4.18, 3.90, 0.30, "= AI 看 「特 征」 来 认 东西!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.65, 4.48, 3.90, 0.25, "AI looks at features to identify things",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

# MIDDLE: big arrow
arrow(s, 4.90, 2.65, w=0.50, h=0.50, color=DAY)
tb(s, 4.85, 3.25, 0.65, 0.30, "现在", sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)

# RIGHT: now we do
panel(s, 5.55, 0.95, 4.05, 3.95, DAY, fill=WHITE, lw=3)
panel_head(s, 5.55, 0.95, 4.05, DAY, "🤖 现在 — 我们 来 做!", sz=12)
tb(s, 5.65, 1.65, 3.85, 0.95, "🤖",
   sz=68, a=PP_ALIGN.CENTER)
tb(s, 5.65, 2.65, 3.85, 0.40, "设 计 AI 机器人!",
   sz=17, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 5.65, 3.05, 3.85, 0.30, "Design our own AI robot!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)
# What this robot does
wr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.70), Inches(3.50), Inches(3.85), Inches(1.30))
wr.fill.solid(); wr.fill.fore_color.rgb = STAR; wr.line.fill.background()
tb(s, 5.80, 3.58, 3.65, 0.32, "🎯 它 会 看 「表情」",
   sz=13, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 5.80, 3.92, 3.65, 0.32, "👀 看 嘴巴 / 眼睛 / 眉毛",
   sz=12, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 5.80, 4.28, 3.65, 0.30, "= 表 情 特 征!",
   sz=12, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 5.80, 4.58, 3.65, 0.22, "= emotion features",
   sz=9, c=DARK, a=PP_ALIGN.LEFT)

# Bottom unifier
unif = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.42))
unif.fill.solid(); unif.fill.fore_color.rgb = DAY; unif.line.fill.background()
tb(s, 0.55, 5.05, 9.0, 0.32, "✨ 我们 学 的 「特 征」 现在 用 在 机器人 上!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 分钟 过 渡:\n• 老师 提问 复习: 「AI 看 什么 来 认 猫?」\n• 学生 回 答: 尖 耳朵, 胡须, 尾巴 ...\n• 引导: 「这些 都 叫 特 征」\n• 然后 切 入: 「现在 我们 让 机器人 看 一 个 不同 的 东 西 — 表 情!」\n• 嘴巴 / 眼睛 / 眉毛 都 是 「表 情 特 征」")


# ============================================================
# PROJECT SLIDE 2 · 项 目 介绍 (AI 表情 识别 机器人)
# ============================================================
s = ns(prs); bg(s, INK, prs)
# Stars decoration
for x, y in [(0.4, 0.45), (9.1, 0.5), (0.5, 4.85), (9.0, 4.85), (1.2, 0.55), (8.3, 4.95)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.40, 9.4, 0.40, "🏆 Project Time · 项 目 时 间!",
   sz=16, b=True, c=NEON, a=PP_ALIGN.CENTER)

# Big project title box
tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.95), Inches(8.8), Inches(1.70))
tt.fill.solid(); tt.fill.fore_color.rgb = DAY
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.15, 8.4, 0.70, "🤖 AI 表情 识别 机器人",
   sz=34, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 1.90, 8.4, 0.40, "Emotion Recognition Robot",
   sz=16, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.30, 8.4, 0.30, "它 会 看 表 情, 猜 心 情!",
   sz=12, c=WARM, a=PP_ALIGN.CENTER)

# Big question
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.85), Inches(8.8), Inches(0.60))
qb.fill.solid(); qb.fill.fore_color.rgb = STAR; qb.line.fill.background()
tb(s, 0.7, 2.92, 8.6, 0.45, "🤔 你 的 机器人 会 认 出 谁 的 心 情?",
   sz=18, b=True, c=INK, a=PP_ALIGN.CENTER)

# Who it can help — 5 examples
helps = [
    ("🧒", "小 朋友"),
    ("👨‍🏫", "老师"),
    ("👨‍👩‍👧", "家人"),
    ("🏥", "病人"),
    ("👫", "朋友"),
]
hw = 1.65; hgap = 0.13
htotal = 5*hw + 4*hgap; hstart = (10 - htotal)/2
for i, (em, txt) in enumerate(helps):
    x = hstart + i*(hw + hgap)
    panel(s, x, 3.65, hw, 1.20, WARM, fill=WHITE, lw=2)
    tb(s, x, 3.75, hw, 0.60, em, sz=34, a=PP_ALIGN.CENTER)
    tb(s, x, 4.40, hw, 0.40, txt, sz=13, b=True, c=DAY, a=PP_ALIGN.CENTER)

# Bottom message
tb(s, 0.3, 5.00, 9.4, 0.32, "💡 你 的 机器人 可以 帮 助 别 人!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟 介绍:\n• 大 声 宣布 项 目 名字 — 让 学生 跟 读 「AI 表情 识别 机器人!」\n• 问 学生: 你 的 机器人 想 帮 谁?\n• 让 学生 想 一 想 自 己 的 故事\n• 例子: 病人 (机器人 知道 你 难过 就 安慰); 老师 (帮 老师 看 全 班 心情)\n• 鼓励 创意 — 任何 答案 都 好!")


# ============================================================
# PROJECT SLIDE 3 · 复习 — AI 看 什么 特 征?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "👀 AI 看 什么? · What Does AI Look At?", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "在 「脸」 上 找 特 征 — 嘴巴 + 眼睛 + 眉毛!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: big face with labels
panel(s, 0.40, 1.20, 4.40, 3.85, DAY, fill=WHITE, lw=3)
# Big face emoji
tb(s, 0.40, 1.40, 4.40, 1.60, "😊", sz=110, a=PP_ALIGN.CENTER)

# Label connectors (using arrows)
labels = [
    ("👁️", "眼 睛", "yǎn jīng", "Eyes", CYBER),
    ("〰️", "眉 毛", "méi máo", "Eyebrows", PURPLE),
    ("👄", "嘴 巴", "zuǐ ba", "Mouth", PINK),
]
for i, (em, cn, py, en, cl) in enumerate(labels):
    y = 3.20 + i*0.55
    tb(s, 0.55, y, 0.50, 0.45, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, 1.10, y+0.02, 1.60, 0.32, cn, sz=15, b=True, c=cl)
    tb(s, 2.75, y+0.05, 1.80, 0.28, f"{py} · {en}", sz=10, c=GRAY)

# RIGHT: 4 emotions with features
panel(s, 5.00, 1.20, 4.60, 3.85, ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.00, 1.20, 4.60, ORANGE, "🎭 4 种 心 情 · 4 Emotions", sz=12)

emotions = [
    ("😊", "开 心", "嘴 弯 上 · 眼 弯", GREEN),
    ("😡", "生 气", "眉 皱 · 嘴 直", RED),
    ("😢", "难 过", "嘴 弯 下 · 眼 垂", CYBER),
    ("😲", "惊 讶", "眼 大 · 嘴 张", PURPLE),
]
for i, (em, cn, feat, cl) in enumerate(emotions):
    y = 1.85 + i*0.75
    tb(s, 5.15, y, 0.60, 0.60, em, sz=32, a=PP_ALIGN.LEFT)
    tb(s, 5.85, y+0.05, 1.25, 0.32, cn, sz=15, b=True, c=cl)
    tb(s, 7.20, y+0.10, 2.30, 0.45, feat, sz=10, b=True, c=DARK)

# Bottom prompt
tb(s, 0.4, 5.10, 9.2, 0.30, "🤔 你 开 心 的 时候 嘴 是 什么 样? 演 一 演!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 体 验:\n• 让 学生 跟 着 做 表 情 (开 心 → 生 气 → 难 过 → 惊 讶)\n• 问: 你 嘴 是 什么 样? 眼 睛 是 什么 样?\n• 老师 在 白 板 写 学 生 答 案\n• 引导: 这 些 就 是 「表 情 特 征」 — 跟 AI 认 猫 一 样\n• K-2 演 表 情, 3-5 可以 解 释 「特 征」 概念")


# ============================================================
# PROJECT SLIDE 4 · 材 料 准备
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🛠️ 材 料 准备 · Materials", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "拿 一 拿 — 准备 好 这 些 就 可以 开 始!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Required materials (top row)
tb(s, 0.4, 1.25, 9.2, 0.28, "✅ 必 备 · Required:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)

required = [
    ("📄", "卡 纸 模板", "Paper template"),
    ("🖍️", "彩 笔", "Markers"),
    ("✂️", "剪刀", "Scissors"),
    ("📎", "胶 水", "Glue"),
]
rw = 2.20; rgap = 0.10
rtotal = 4*rw + 3*rgap; rstart = (10 - rtotal)/2
for i, (em, cn, en) in enumerate(required):
    x = rstart + i*(rw + rgap)
    panel(s, x, 1.65, rw, 1.40, DAY, fill=WHITE, lw=2.5)
    tb(s, x, 1.78, rw, 0.65, em, sz=34, a=PP_ALIGN.CENTER)
    tb(s, x, 2.45, rw, 0.32, cn, sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.78, rw-0.10, 0.25, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Optional materials (bottom row)
tb(s, 0.4, 3.25, 9.2, 0.28, "⭐ 可 选 · Optional (更 好 看!):",
   sz=11, b=True, c=ORANGE, a=PP_ALIGN.LEFT)

optional = [
    ("📍", "大 头 针", "Paper fasteners"),
    ("🔘", "魔术 贴", "Velcro"),
    ("👀", "活动 眼睛", "Googly eyes"),
    ("✨", "贴 纸 / 锡 纸", "Stickers / foil"),
]
for i, (em, cn, en) in enumerate(optional):
    x = rstart + i*(rw + rgap)
    panel(s, x, 3.60, rw, 1.30, ORANGE, fill=WARM, lw=2)
    tb(s, x, 3.72, rw, 0.55, em, sz=28, a=PP_ALIGN.CENTER)
    tb(s, x, 4.28, rw, 0.30, cn, sz=11, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 4.58, rw-0.10, 0.25, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom tip
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
tip.fill.solid(); tip.fill.fore_color.rgb = DAY; tip.line.fill.background()
tb(s, 0.55, 5.10, 9.0, 0.30, "👉 大头 针 / 魔术 贴 让 表情 可以 换! Use fasteners to swap emotions!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟 发 材料:\n• 必 备 材料 都 是 教 室 里 常 见 的\n• 可 选 材料 让 项目 更 有 趣 — 但 不 是 必 须 的\n• 重 点: 大头 针 / 魔术 贴 让 「表 情」 可以 换 — 这 样 一 个 机器人 可以 显 示 多 种 表 情\n• 如 果 没 有 大头 针: 可以 用 双 面 胶 临 时 贴 上 / 撕 下")


# ============================================================
# PROJECT SLIDE 5 · 第 一 步 · 做 机器人 脸
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "1️⃣ 第 一 步 · 做 机器人 脸 · Make Robot Face", CYBER)

tb(s, 0.4, 0.85, 9.2, 0.30, "先 设 计 你 自己 的 机器人 — 用 你 的 创 意!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: example robot face
panel(s, 0.40, 1.20, 4.40, 3.85, CYBER, fill=WHITE, lw=3)
tb(s, 0.40, 1.30, 4.40, 0.40, "👇 例子 · Example",
   sz=11, b=True, c=CYBER, a=PP_ALIGN.CENTER)
# Big robot
tb(s, 0.40, 1.75, 4.40, 1.85, "🤖",
   sz=120, a=PP_ALIGN.CENTER)
# Name plate
np_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.80), Inches(3.85), Inches(3.60), Inches(0.85))
np_box.fill.solid(); np_box.fill.fore_color.rgb = STAR; np_box.line.fill.background()
tb(s, 0.85, 3.90, 3.50, 0.30, "📛 我 的 名字:",
   sz=11, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 0.85, 4.20, 3.50, 0.50, "____________",
   sz=22, b=True, c=INK, a=PP_ALIGN.CENTER)

# RIGHT: 3-step process
panel(s, 5.00, 1.20, 4.60, 3.85, ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.00, 1.20, 4.60, ORANGE, "📋 怎么 做 · How To", sz=12)

face_steps = [
    ("🖍️", "涂 色", "Color", "用 你 喜 欢 的 颜色"),
    ("✨", "装 饰", "Decorate", "贴 纸 / 锡 纸 / 眼 睛"),
    ("✂️", "剪 出 来", "Cut out", "小 心 剪 — 慢 慢 来"),
]
for i, (em, cn, en, hint) in enumerate(face_steps):
    y = 1.85 + i*0.75
    tb(s, 5.15, y, 0.55, 0.55, em, sz=26, a=PP_ALIGN.LEFT)
    tb(s, 5.80, y+0.05, 1.50, 0.30, cn, sz=14, b=True, c=ORANGE)
    tb(s, 5.80, y+0.35, 1.50, 0.25, en, sz=9, c=GRAY)
    tb(s, 7.40, y+0.10, 2.15, 0.50, hint, sz=10, b=True, c=DARK)

# Bottom prompt
prompt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
prompt.fill.solid(); prompt.fill.fore_color.rgb = CYBER; prompt.line.fill.background()
tb(s, 0.55, 5.10, 9.0, 0.30, "🤔 它 叫 什么 名字? 想 一 个 酷 的!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 第 一 步:\n• 强调: 不 用 画 得 完 美 — 创 意 最 重 要!\n• 鼓励 学生 给 机器人 取 中文 名字\n• 例子 名字: 智 智 / AI 宝宝 / 表 情 王 / 心 情 队 长\n• K-2: 老师 帮 写 名字\n• 3-5: 自己 想 + 写")


# ============================================================
# PROJECT SLIDE 6 · 第 二 步 · 做 不同 表 情
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "2️⃣ 第 二 步 · 做 不同 表 情 · Make Emotions", PURPLE)

tb(s, 0.4, 0.85, 9.2, 0.30, "AI 要 学 不 同 表 情 — 你 的 机器人 也 要!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "AI needs to learn different emotions — your robot does too!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 4 emotion cards (interchangeable mouths/faces)
emotions = [
    ("😊", "开 心", "Happy", "嘴 弯 上 ⌣", GREEN),
    ("😡", "生 气", "Angry", "眉 皱 + 嘴 直", RED),
    ("😢", "难 过", "Sad", "嘴 弯 下 ⌢", CYBER),
    ("😲", "惊 讶", "Surprised", "嘴 圆 O", PURPLE),
]
ew = 2.15; egap = 0.15
etotal = 4*ew + 3*egap; estart = (10 - etotal)/2
for i, (em, cn, en, feat, cl) in enumerate(emotions):
    x = estart + i*(ew + egap)
    panel(s, x, 1.65, ew, 2.85, cl, fill=WHITE, lw=3)
    tb(s, x, 1.80, ew, 1.00, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 2.85, ew, 0.40, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 3.25, ew, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Feature hint
    fh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.10), Inches(3.65), Inches(ew-0.20), Inches(0.65))
    fh.fill.solid(); fh.fill.fore_color.rgb = WARM; fh.line.fill.background()
    tb(s, x+0.10, 3.72, ew-0.20, 0.25, "🎯 特 征:", sz=9, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.97, ew-0.20, 0.30, feat, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Bottom tip
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.80))
tip.fill.solid(); tip.fill.fore_color.rgb = PURPLE
tip.line.color.rgb = STAR; tip.line.width = Pt(2)
tb(s, 0.55, 4.75, 9.0, 0.35, "💡 用 大头 针 / 魔术 贴 — 可以 换 表 情!",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.28, "Use paper fasteners or velcro so your robot can swap emotions!",
   sz=10, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10-15 分钟 第 二 步:\n• 学生 做 4 张 嘴巴 / 眉毛 卡片 — 每 张 代 表 一 种 表 情\n• K-2 老师 准备 模板, 学生 涂 色 + 剪\n• 3-5 自己 设 计 + 剪\n• 关 键: 表 情 要 「清 楚」 — 让 别 人 看 出 来\n• 这 是 训练 数据 的 类 比: 不 清 楚 的 训练 数据 = AI 学 不 好")


# ============================================================
# PROJECT SLIDE 7 · 第 三 步 · 让 机器人 会 「认」
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "3️⃣ 第 三 步 · 让 机器人 会 「认」 · Teach It!", PINK)

tb(s, 0.4, 0.85, 9.2, 0.30, "告诉 你 的 机器人: 它 看 什么 「特 征」 来 猜 心 情?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Big sentence frames panel
sf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.30), Inches(9.20), Inches(2.75))
sf.fill.solid(); sf.fill.fore_color.rgb = INK
sf.line.color.rgb = STAR; sf.line.width = Pt(3)
tb(s, 0.55, 1.40, 9.0, 0.35, "💬 写 一 写 — 你 的 机器人 规则:",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)

# 3 sentence frames
frames = [
    ("👀", "我 的 机器人 看 ______ 。", "(嘴巴 / 眼睛 / 眉毛)"),
    ("😊", "它 看 到 ______ 就 觉得 你 开 心。", "(嘴 弯 上 / 微笑 ...)"),
    ("😡", "它 看 到 ______ 就 觉得 你 生 气。", "(眉 皱 / 嘴 紧 ...)"),
]
for i, (em, frame, hint) in enumerate(frames):
    y = 1.90 + i*0.70
    tb(s, 0.65, y, 0.55, 0.45, em, sz=24, a=PP_ALIGN.LEFT)
    tb(s, 1.25, y+0.05, 8.20, 0.35, frame,
       sz=15, b=True, c=STAR, a=PP_ALIGN.LEFT)
    tb(s, 1.25, y+0.42, 8.20, 0.25, hint,
       sz=9, c=WARM, a=PP_ALIGN.LEFT)

# Connection callout
con = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.20), Inches(9.20), Inches(0.85))
con.fill.solid(); con.fill.fore_color.rgb = PINK
con.line.color.rgb = STAR; con.line.width = Pt(2.5)
tb(s, 0.55, 4.30, 9.0, 0.35, "🔗 这 些 就 是 你 的 机器人 「特 征 表」!",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.70, 9.0, 0.28, "These are your robot's 'features list' — just like a real AI!",
   sz=10, c=STAR, a=PP_ALIGN.CENTER)

tb(s, 0.4, 5.15, 9.2, 0.28, "✏️ K-2 老师 帮 写; 3-5 自己 写!",
   sz=11, b=True, c=PINK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 第 三 步 (核心 连接 ML 概念!):\n• 学生 写 句 型 — 定义 自己 机器人 的 规则\n• 例子:\n  - 我 的 机器人 看 嘴巴\n  - 它 看 到 嘴 弯 上 就 觉得 你 开 心\n  - 它 看 到 嘴 弯 下 就 觉得 你 难过\n• K-2 老师 帮 写; 3-5 自己 写\n• 这 是 项目 跟 Session 1 ML 概念 最 强 的 连接 点\n• 让 学生 意识 到: 自己 也 在 「写 规则」 训练 AI")


# ============================================================
# PROJECT SLIDE 8 · 测试 你 的 机器人!
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🧪 测试 你 的 机器人! · Test Your Robot!", ORANGE)

tb(s, 0.4, 0.85, 9.2, 0.30, "找 同 桌 一 起 玩 — 看 你 的 机器人 准 不 准!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Partner activity flow (3 steps)
panel(s, 0.40, 1.25, 9.20, 1.95, ORANGE, fill=WHITE, lw=3)
panel_head(s, 0.40, 1.25, 9.20, ORANGE, "👫 配 对 游戏 · Partner Activity", sz=12)

partner_steps = [
    ("A", "同学 A 做 一 个 表 情",
     "Student A makes a face", GREEN),
    ("B", "同学 B 用 机器人 猜!",
     "Student B uses robot to guess", CYBER),
    ("🔄", "换 一 换 — 再 玩!",
     "Switch + play again!", PINK),
]
sw = 2.85; sgap = 0.20
stotal = 3*sw + 2*sgap; sstart = (10 - stotal)/2
for i, (badge, cn, en, cl) in enumerate(partner_steps):
    x = sstart + i*(sw + sgap)
    # Mini card inside panel
    mc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85), Inches(sw), Inches(1.30))
    mc.fill.solid(); mc.fill.fore_color.rgb = WARM
    mc.line.color.rgb = cl; mc.line.width = Pt(2)
    # Badge circle
    bc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+sw/2-0.30), Inches(1.95), Inches(0.60), Inches(0.60))
    bc.fill.solid(); bc.fill.fore_color.rgb = cl; bc.line.fill.background()
    tb(s, x+sw/2-0.30, 2.00, 0.60, 0.50, badge, sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x, 2.62, sw, 0.30, cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.92, sw-0.10, 0.22, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Reflection questions
panel(s, 0.40, 3.40, 9.20, 1.55, DAY, fill=WARM, lw=3)
panel_head(s, 0.40, 3.40, 9.20, DAY, "🤔 一 起 讨 论  Discuss Together", sz=12)
questions = [
    ("✅", "猜 对 了 吗?", "Was it right?"),
    ("❓", "为什么?", "Why or why not?"),
    ("🔍", "哪 里 看 错 了?", "Where did it go wrong?"),
]
qw = 2.85; qgap = 0.20
qtotal = 3*qw + 2*qgap; qstart = (10 - qtotal)/2
for i, (em, cn, en) in enumerate(questions):
    x = qstart + i*(qw + qgap)
    tb(s, x, 3.95, qw, 0.45, em, sz=22, a=PP_ALIGN.CENTER)
    tb(s, x, 4.40, qw, 0.30, cn, sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x, 4.70, qw, 0.22, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom connection
con = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
con.fill.solid(); con.fill.fore_color.rgb = DAY; con.line.fill.background()
tb(s, 0.55, 5.10, 9.0, 0.30, "🤖 像 真 正 训练 AI 一 样 — 我们 也 要 测试!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 测试 活 动:\n• 学生 找 同 桌 配 对\n• 一 个 做 表 情, 一 个 「用」 机器人 看 嘴巴/眼睛/眉毛 → 看 是 哪 种 心情\n• 然后 互 换\n• 每 组 测 3-5 次\n• 老师 引导: 「这 就 是 真 AI 训练 的 过程 — 测试 才 知 道 行 不 行!」\n• 让 学 生 自然 进入 下 一 张 的 「改 进」 讨论")


# ============================================================
# PROJECT SLIDE 9 · 机器人 犯 错 怎么 办?
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔧 机器人 犯 错 怎么 办? · What If It's Wrong?", PINK)

tb(s, 0.4, 0.85, 9.2, 0.30, "AI 工程师 不 怕 错 — 错 了 就 改!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "AI engineers don't fear mistakes — they fix them!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 3 improvement ideas
improvements = [
    ("✏️", "把 特 征 画 更 清 楚",
     "Make features clearer",
     "嘴巴 / 眼睛 画 得 更 大 更 清", CYBER),
    ("➕", "做 更 多 表 情",
     "Add more emotion examples",
     "增 加 5 种 表 情 → 数据 更 多", ORANGE),
    ("📝", "改 规 则",
     "Change the rules",
     "重新 想: 它 应 该 看 什么?", PURPLE),
]
iw = 2.95; igap = 0.18
itotal = 3*iw + 2*igap; istart = (10 - itotal)/2
for i, (em, cn, en, hint, cl) in enumerate(improvements):
    x = istart + i*(iw + igap)
    panel(s, x, 1.65, iw, 2.85, cl, fill=WHITE, lw=3)
    tb(s, x, 1.78, iw, 0.85, em, sz=46, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, iw, 0.45, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x, 3.18, iw, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Hint box
    hb_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.15), Inches(3.60), Inches(iw-0.30), Inches(0.80))
    hb_box.fill.solid(); hb_box.fill.fore_color.rgb = WARM; hb_box.line.fill.background()
    tb(s, x+0.15, 3.70, iw-0.30, 0.55, hint, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Connection to ML
con = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.70), Inches(9.20), Inches(0.75))
con.fill.solid(); con.fill.fore_color.rgb = DAY
con.line.color.rgb = STAR; con.line.width = Pt(2.5)
tb(s, 0.55, 4.80, 9.0, 0.35, "💡 这 就 是 机器 学习 改 进 的 方法!",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.18, 9.0, 0.25, "This is exactly how machine learning improves!",
   sz=10, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 改 进 讨论:\n• 测 完 后, 让 学 生 想: 怎么 让 机器人 更 准?\n• 关键 联 系: 真 AI 也 是 这 样 改 进\n  - 更 多 训练 数据 = 我们 做 更 多 表 情\n  - 更 清 楚 特 征 = 画 得 更 清\n  - 改 规 则 = 重新 训练\n• 让 学 生 试 着 改 进 自己 的 机器人 (5 分钟)")


# ============================================================
# PROJECT SLIDE 10 · 升 级 挑 战 (3-5 高 年 级)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "⭐ 升 级 挑战 · Extra Challenge (3-5)", PURPLE)

tb(s, 0.4, 0.85, 9.2, 0.30, "3-5 年 级 — 让 你 的 机器人 认 更 难 的 表 情!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 harder emotions
harder = [
    ("😕", "困 惑", "Confused"),
    ("😴", "累", "Tired"),
    ("😳", "害 羞", "Embarrassed"),
    ("🤩", "兴 奋", "Excited"),
]
hw = 2.15; hgap = 0.15
htotal = 4*hw + 3*hgap; hstart = (10 - htotal)/2
for i, (em, cn, en) in enumerate(harder):
    x = hstart + i*(hw + hgap)
    panel(s, x, 1.30, hw, 2.10, PURPLE, fill=WHITE, lw=3)
    tb(s, x, 1.45, hw, 0.95, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 2.50, hw, 0.40, cn, sz=17, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.92, hw-0.10, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Big question
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.65), Inches(9.20), Inches(1.10))
qb.fill.solid(); qb.fill.fore_color.rgb = PURPLE
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 3.78, 9.0, 0.45, "🤔 这 些 更 难 吗? 为 什么?",
   sz=20, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.25, 9.0, 0.30, "Are these harder? Why?",
   sz=12, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.55, 9.0, 0.22, "💡 提 示: 表 情 有 时 很 像 — 困 惑 vs 难 过?",
   sz=9, b=True, c=WARM, a=PP_ALIGN.CENTER)

# Bottom guide
tb(s, 0.4, 4.95, 9.2, 0.30, "🎯 高 年 级 — 想 想 AI 怎么 区 分 「累」 vs 「难 过」?",
   sz=12, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 (高 年 级):\n• K-2 可 跳 过 — 继续 玩 4 个 基本 表 情\n• 3-5 加 这 4 个 表 情\n• 关 键 讨论:\n  - 困 惑 vs 难 过: 都 是 嘴 弯 下 — 怎么 区 分?\n  - 累 vs 难 过: 眼 都 半 闭 — 怎么 区 分?\n• 引 导: AI 也 有 「类 似 类 别 难 区 分」 的 问 题\n• 这 是 高 年 级 的 critical thinking 时 刻")


# ============================================================
# PROJECT SLIDE 11 · 反 思
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🤔 今天 你 像 AI 工程师 吗? · Are You an AI Engineer?", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "做 完 项目 后 — 一 起 想 一 想!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: 3 reflection questions
panel(s, 0.40, 1.25, 4.55, 3.75, CYBER, fill=WHITE, lw=3)
panel_head(s, 0.40, 1.25, 4.55, CYBER, "❓ 想 一 想  Reflect", sz=12)
ref_qs = [
    ("🧠", "AI 怎么 学 习?"),
    ("⚠️", "为什么 会 犯 错?"),
    ("👀", "你 的 机器人 看 什么 特 征?"),
]
for i, (em, q) in enumerate(ref_qs):
    y = 1.85 + i*0.95
    tb(s, 0.55, y, 0.55, 0.55, em, sz=28, a=PP_ALIGN.LEFT)
    tb(s, 1.15, y+0.12, 3.65, 0.65, q, sz=13, b=True, c=DARK)

# RIGHT: sentence frames
panel(s, 5.05, 1.25, 4.55, 3.75, ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.05, 1.25, 4.55, ORANGE, "💬 句 型  Sentence Frames", sz=12)
ref_frames = [
    "AI 需要 ______ 。",
    "AI 看 ______ 。",
    "我 的 机器人 帮助 ______ 。",
]
for i, frame in enumerate(ref_frames):
    y = 1.85 + i*0.95
    # Frame card
    fc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.20), Inches(y), Inches(4.25), Inches(0.75))
    fc.fill.solid(); fc.fill.fore_color.rgb = WARM; fc.line.fill.background()
    tb(s, 5.30, y+0.05, 4.05, 0.30, f"{i+1}.", sz=11, b=True, c=ORANGE, a=PP_ALIGN.LEFT)
    tb(s, 5.30, y+0.32, 4.05, 0.42, frame, sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Bottom encouragement
enc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.10), Inches(9.20), Inches(0.32))
enc.fill.solid(); enc.fill.fore_color.rgb = DAY; enc.line.fill.background()
tb(s, 0.55, 5.13, 9.0, 0.25, "🌟 你 今天 是 AI 工程师!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 反 思:\n• 让 学 生 用 句 型 说 — 不 只 是 想\n• 鼓励 K-2 用 简 单 中 文 答\n• 3-5 可以 写 在 纸 上\n• 关 键 答 案:\n  - AI 需要 「数据 / 例子 / 训练 / 标签」\n  - AI 看 「特 征」\n  - 我 的 机器人 帮助 「老师 / 病人 / 家人 ...」")


# ============================================================
# PROJECT SLIDE 12 · 展览 — Gallery Walk
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎤 AI Robot 展览! · Gallery Walk", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "拿 你 的 机器人 — 一起 走 一 走 看 看!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "Bring your robot — let's walk around and see everyone's!",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 4 sentence frames for presentation
panel(s, 0.40, 1.55, 9.20, 3.05, DAY, fill=INK, lw=3)
tb(s, 0.55, 1.70, 9.0, 0.35, "💬 介绍 你 的 机器人 — 用 这 些 句 型:",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)

show_frames = [
    ("📛", "我 的 机器人 叫 ______ 。"),
    ("💖", "它 帮助 ______ 。"),
    ("👀", "它 看 ______ 。"),
    ("✨", "它 可以 认 出 ______ 。"),
]
for i, (em, frame) in enumerate(show_frames):
    y = 2.15 + i*0.58
    tb(s, 0.75, y, 0.50, 0.45, em, sz=24, a=PP_ALIGN.LEFT)
    tb(s, 1.30, y+0.08, 8.20, 0.45, frame,
       sz=17, b=True, c=STAR, a=PP_ALIGN.LEFT)

# Walking flow + cheers
walk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.70))
walk.fill.solid(); walk.fill.fore_color.rgb = STAR
walk.line.color.rgb = DAY; walk.line.width = Pt(2)
tb(s, 0.55, 4.85, 9.0, 0.32, "🚶 走 一 走! 听 同学 介绍 — 给 掌 声!",
   sz=14, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.18, 9.0, 0.22, "Walk around, listen, applaud each other!",
   sz=9, c=DARK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "10 分钟 展览:\n• 把 桌 子 围 成 圈 / 学生 站 在 自己 位 置 旁\n• 老师 给 顺 序 — 每 个 学 生 30 秒 介绍\n• 听 完 给 掌 声!\n• 鼓励 用 中文 句 型 — 别 怕 错\n• K-2 老师 帮 念; 3-5 自己 说\n• 拍 照 留 念 — 发 给 家 长!\n\n备 案: 如 学 生 太 多, 分 小 组 内 部 展 览 (4-5 人 / 组)")


# ============================================================
# TM PROJECT FOR 3-5 (14 slides) — Teachable Machine guided training
# Older students work in groups of 2-3 on school computers.
# Younger students do the robot craft above; older students do this.
# ============================================================


# ============================================================
# TM-1 · TITLE — 高 年 级 挑 战
# ============================================================
s = ns(prs); bg(s, INK, prs)
# Sparkle decoration
for x, y in [(0.4, 0.45), (9.1, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()

tb(s, 0.3, 0.40, 9.4, 0.40, "🚀 升 级 版 · For Grades 3-5",
   sz=16, b=True, c=NEON, a=PP_ALIGN.CENTER)

# Big title box
tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.00), Inches(8.8), Inches(2.10))
tt.fill.solid(); tt.fill.fore_color.rgb = DAY
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.20, 8.4, 0.85, "高 年 级 挑 战",
   sz=44, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.10, 8.4, 0.50, "训练 自己 的 AI 模型!",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.65, 8.4, 0.35, "Train your own AI model!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

# Setup info banner
si = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.30), Inches(8.8), Inches(0.80))
si.fill.solid(); si.fill.fore_color.rgb = WHITE
si.line.color.rgb = STAR; si.line.width = Pt(2)
tb(s, 0.75, 3.45, 8.5, 0.50, "👥 2-3 人 一 组  ·  💻 学校 电脑  ·  🌐 Teachable Machine  ·  ⏱️ 30 分钟",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)

# Visual emoji row
tb(s, 0.3, 4.25, 9.4, 0.85, "👨‍💻   👩‍💻   💻   🤖   ✨",
   sz=46, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "1 分钟 介 绍:\n• 大 声 宣 布: 「现在 是 高 年 级 挑战 时间!」\n• 重 点 不是 「做 得 完美」, 而是 「观察 AI 怎么 学 + 怎么 错 + 怎么 改」\n• 老师 提前 准备: 笔记本 + 摄像 头 + 已 打开 Teachable Machine\n• K-2 学生 同 时 在 做 robot craft (parallel track)")


# ============================================================
# TM-2 · GROUP RULES — 小 组 规 则
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "📋 小 组 规 则 · Group Rules", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "开 始 前 — 请 先 看 这 6 条 规 则!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 6 rules in 2 rows of 3
rules = [
    ("1️⃣", "👥", "每 组 2-3 人", "Groups of 2-3"),
    ("2️⃣", "🎯", "只 选 一 个 主 题", "Pick ONE topic"),
    ("3️⃣", "🔢", "只 做 两 个 类 别", "Only TWO classes"),
    ("4️⃣", "📷", "每 类 至少 拍 20 张", "20+ photos per class"),
    ("5️⃣", "🧪", "训 练 后 一定 要 测 试", "Always test after training"),
    ("6️⃣", "🔧", "发现 错误 → 改 进", "Find errors → improve"),
]
rw = 2.95; rgap_x = 0.18; rgap_y = 0.20; rh = 1.55
rstart_x = (10 - 3*rw - 2*rgap_x)/2
for i, (num, em, cn, en) in enumerate(rules):
    row = i // 3; col = i % 3
    x = rstart_x + col*(rw + rgap_x)
    y = 1.30 + row*(rh + rgap_y)
    panel(s, x, y, rw, rh, DAY, fill=WHITE, lw=2.5)
    # Number badge top-left
    tb(s, x+0.10, y+0.08, 0.50, 0.40, num, sz=18, b=True, c=DAY, a=PP_ALIGN.LEFT)
    # Emoji center
    tb(s, x, y+0.10, rw, 0.55, em, sz=28, a=PP_ALIGN.CENTER)
    # CN text
    tb(s, x, y+0.78, rw, 0.40, cn, sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
    # EN subtitle
    tb(s, x+0.05, y+1.18, rw-0.10, 0.28, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom teacher note
tn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
tn.fill.solid(); tn.fill.fore_color.rgb = STAR; tn.line.fill.background()
tb(s, 0.55, 5.10, 9.0, 0.30, "💡 保 持 简 单 — 只 做 2 个 类 别 就 好!",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟 看 规 则:\n• 老师 大 声 念 一 遍\n• 学 生 一 起 跟 着 读\n• 强 调 第 3 条: 只 做 2 个 类 别 — 不 要 贪 多\n• 强 调 第 4 条: 20+ 张 是 「最 少」, 越 多 越 好")


# ============================================================
# TM-3 · ML PROCESS — 机器 学习 六 步 (consistent with slides 6 + 23)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔄 机器 学习 六 步 · The 6 Steps", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "做 项 目 时 — 跟 着 早 上 学 的 6 步!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Same 6 steps as overview slide 6 + wrap-up slide 23 — KEEP CONSISTENT
process_steps = [
    ("1️⃣", "📷", "收集 例子", "Collect", CYBER),
    ("2️⃣", "🏷️", "贴 标签", "Label", ORANGE),
    ("3️⃣", "🧠", "AI 学 习", "Study", PURPLE),
    ("4️⃣", "🔍", "找 规律", "Patterns", PINK),
    ("5️⃣", "🧪", "测 试", "Test", DAY),
    ("6️⃣", "✨", "做 判 断", "Predict", GREEN),
]
# 2 rows of 3 layout (matches slide 6 overview style)
sw = 2.85; sgap_x = 0.20; sgap_y = 0.20; sh = 1.50
sstart_x = (10 - 3*sw - 2*sgap_x)/2
for i, (num, em, cn, en, cl) in enumerate(process_steps):
    row = i // 3; col = i % 3
    x = sstart_x + col*(sw + sgap_x)
    y = 1.30 + row*(sh + sgap_y)
    panel(s, x, y, sw, sh, cl, fill=WHITE, lw=2.5)
    # Number badge top-left
    tb(s, x+0.10, y+0.08, 0.55, 0.45, num, sz=18, b=True, c=cl, a=PP_ALIGN.LEFT)
    # Big emoji centered
    tb(s, x+0.60, y+0.08, sw-0.70, 0.70, em, sz=36, a=PP_ALIGN.CENTER)
    # Chinese name
    tb(s, x, y+0.82, sw, 0.38, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    # English subtitle
    tb(s, x, y+1.18, sw, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Arrow to next (same row)
    if col < 2:
        arrow(s, x + sw + 0.02, y + sh/2 - 0.12, w=0.14, h=0.24, color=DAY)

# Big reminder
rem = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.50))
rem.fill.solid(); rem.fill.fore_color.rgb = DAY
rem.line.color.rgb = STAR; rem.line.width = Pt(2.5)
tb(s, 0.55, 5.05, 9.0, 0.32, "💡 数据 越 清 楚, AI 越 容易 学 会!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 跟 早 上 的 6 步 一 样 — 帮 学 生 加 深 印象\n• 一起 数 「1, 2, 3, 4, 5, 6 步!」\n• 在 Teachable Machine 项目 里, 学 生 主要 做 步 1+2+5 (收集 + 贴 标签 + 测 试)\n• AI 自动 做 步 3+4+6 (学 习 / 找 规律 / 做 判 断)")


# ============================================================
# TM-4 · WEBSITE STEPS — 打 开 Teachable Machine
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🌐 打 开 Teachable Machine · How to Open", DAY)

# URL banner at top
ub = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.85), Inches(9.20), Inches(0.42))
ub.fill.solid(); ub.fill.fore_color.rgb = INK; ub.line.fill.background()
ub.click_action.hyperlink.address = "https://teachablemachine.withgoogle.com/"
url_tb = s.shapes.add_textbox(Inches(0.40), Inches(0.92), Inches(9.20), Inches(0.30))
url_tf = url_tb.text_frame
url_p = url_tf.paragraphs[0]; url_p.alignment = PP_ALIGN.CENTER
url_run = url_p.add_run()
url_run.text = "🔗 teachablemachine.withgoogle.com  (点 这 里 打 开)"
url_run.font.size = Pt(13); url_run.font.bold = True
url_run.font.color.rgb = STAR; url_run.font.name = 'KaiTi'
url_run.hyperlink.address = "https://teachablemachine.withgoogle.com/"

# 8 steps in 2 columns
website_steps = [
    ("1.", "打 开 网址", "Go to teachablemachine.withgoogle.com"),
    ("2.", "点 「Get Started」", "Click 'Get Started'"),
    ("3.", "选 「Image Project」", "Choose 'Image Project'"),
    ("4.", "选 「Standard Image Model」", "Choose 'Standard Image Model'"),
    ("5.", "重 命 名 Class 1 + Class 2", "Rename Class 1 and Class 2"),
    ("6.", "用 摄像 头 收集 图片", "Use webcam to collect images"),
    ("7.", "点 「Train Model」", "Click 'Train Model'"),
    ("8.", "测 试 你 的 模型!", "Test your model!"),
]
sw2 = 4.40; sgap2 = 0.20; sh2 = 0.48
sstart_x2 = (10 - 2*sw2 - sgap2)/2
for i, (num, cn, en) in enumerate(website_steps):
    row = i // 2; col = i % 2
    x = sstart_x2 + col*(sw2 + sgap2)
    y = 1.50 + row*(sh2 + 0.05)
    # Step card
    sc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(sw2), Inches(sh2))
    sc.fill.solid(); sc.fill.fore_color.rgb = WHITE
    sc.line.color.rgb = DAY; sc.line.width = Pt(1.5)
    tb(s, x+0.10, y+0.08, 0.40, 0.35, num, sz=14, b=True, c=DAY, a=PP_ALIGN.LEFT)
    tb(s, x+0.55, y+0.06, sw2-0.65, 0.22, cn, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, x+0.55, y+0.26, sw2-0.65, 0.20, en, sz=8, c=GRAY, a=PP_ALIGN.LEFT)

# Bottom tip
bt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.42))
bt.fill.solid(); bt.fill.fore_color.rgb = STAR; bt.line.fill.background()
tb(s, 0.55, 5.05, 9.0, 0.32, "💡 按 钮 是 英文 的 — 老师 帮 你 找!",
   sz=12, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 分钟:\n• 老师 在 投影 演示 一 遍\n• 学生 跟 着 一 步 一 步 做\n• 英文 按 钮: Get Started / Image Project / Standard Image Model / Train Model\n• 如 摄像 头 不 工 作, 检 查 浏览 器 权 限\n• 准 备 备 案: 如 网络 慢, 可 提前 打 开")


# ============================================================
# TM-5 · CHOOSE ONE PROJECT — 4 options
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎯 选 一 个 训练 任 务 · Pick ONE Project", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "每 组 只 选 一 个 — 不 要 贪 多!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.18, 9.2, 0.26, "Each group picks ONE option only",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# 4 project options in 2x2 grid
project_opts = [
    ("A", "✋", "手势 识别", "Gesture", "石头/剪刀/布", CYBER),
    ("B", "✏️", "文具 识别", "School Supplies", "铅笔/橡皮/书", ORANGE),
    ("C", "🌈", "颜色 物品 分 类", "Color Objects", "红色/蓝色 物品", PINK),
    ("D", "🔷", "形状 识别", "Shapes", "圆形/三角形", PURPLE),
]
pw = 4.40; pgap_x = 0.20; pgap_y = 0.20; ph = 1.55
pstart_x = (10 - 2*pw - pgap_x)/2
for i, (letter, em, cn, en, ex, cl) in enumerate(project_opts):
    row = i // 2; col = i % 2
    x = pstart_x + col*(pw + pgap_x)
    y = 1.55 + row*(ph + pgap_y)
    panel(s, x, y, pw, ph, cl, fill=WHITE, lw=3)
    # Letter badge
    lb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.15), Inches(y+0.20), Inches(0.65), Inches(0.65))
    lb.fill.solid(); lb.fill.fore_color.rgb = cl; lb.line.fill.background()
    tb(s, x+0.15, y+0.28, 0.65, 0.50, letter, sz=22, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Emoji
    tb(s, x+0.95, y+0.12, 1.00, 0.85, em, sz=44, a=PP_ALIGN.LEFT)
    # CN title
    tb(s, x+2.05, y+0.12, pw-2.15, 0.40, cn, sz=17, b=True, c=cl)
    # EN
    tb(s, x+2.05, y+0.55, pw-2.15, 0.28, en, sz=10, c=GRAY)
    # Example
    tb(s, x+0.20, y+1.05, pw-0.40, 0.35, f"例子: {ex}", sz=11, b=True, c=DARK)

# Bottom note
bn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
bn.fill.solid(); bn.fill.fore_color.rgb = DAY; bn.line.fill.background()
tb(s, 0.55, 5.10, 9.0, 0.30, "👉 选 好 后 看 后 面 的 详 细 说 明!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟:\n• 每 组 讨论 30 秒 后 选 一 个\n• 推 荐 难 度:\n  - A 手势 (最 易): 摄像 头 直接 拍, 不 需要 道 具\n  - B 文具 (易): 学生 自带\n  - C 颜色 (中): 需要 多 种 颜色 物品\n  - D 形状 (难 一 些): 画 / 剪 形状\n• 如 多 组 选 同 一 个 — OK, 各 自 做")


# ============================================================
# TM-6 · OPTION A — 手势 识别 模型
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "✋ Option A · 手势 识别 模型 · Gesture AI", CYBER)

# LEFT: 3 gesture pair options
panel(s, 0.40, 0.95, 4.55, 4.10, CYBER, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, CYBER, "🎯 选 一 对 手势  Pick a Pair", sz=12)
pair_opts_a = [
    ("✋ vs ✊", "张开 的 手 vs 拳头"),
    ("👍 vs 👎", "点 赞 vs 倒 赞"),
    ("✌️ vs 🫶", "剪刀 手 vs 比 心"),
]
for i, (visual, cn) in enumerate(pair_opts_a):
    y = 1.55 + i*1.05
    pc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(4.25), Inches(0.95))
    pc.fill.solid(); pc.fill.fore_color.rgb = WARM; pc.line.fill.background()
    tb(s, 0.65, y+0.10, 4.05, 0.45, visual, sz=24, b=True, c=CYBER, a=PP_ALIGN.CENTER)
    tb(s, 0.65, y+0.55, 4.05, 0.32, cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# RIGHT TOP: data collection
panel(s, 5.15, 0.95, 4.45, 2.10, ORANGE, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 0.95, 4.45, ORANGE, "📷 收集 数据 怎么 拍?", sz=11)
data_tips_a = [
    "• 每 类 至少 20 张",
    "• 不 同 同学 的 手",
    "• 不同 角度",
    "• 拍 近 + 清 楚",
]
for i, tip in enumerate(data_tips_a):
    tb(s, 5.30, 1.50 + i*0.35, 4.15, 0.30, tip, sz=10, b=True, c=DARK)

# RIGHT BOT: test + reflection
panel(s, 5.15, 3.20, 4.45, 1.85, PINK, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 3.20, 4.45, PINK, "🧪 测 试 + 反 思", sz=11)
tb(s, 5.30, 3.78, 4.15, 0.28, "🤔 测试 问 题:", sz=10, b=True, c=PINK)
tb(s, 5.30, 4.05, 4.15, 0.25, "• 换 别 人 的 手 — 还 认 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.30, 4.15, 0.25, "• 远 一 点 拍 — 还 准 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.62, 4.15, 0.25, "💡 AI 看 什么? 手指 数? 形状?", sz=10, b=True, c=PINK, a=PP_ALIGN.LEFT)

# Bottom prompt
tb(s, 0.4, 5.15, 9.2, 0.25, "💪 最 简单 的 选项 — 适 合 第 一 次 玩!",
   sz=10, b=True, c=CYBER, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "如 选 A:\n• 最 容易 上手\n• 不 需要 任 何 道 具\n• 注意: 不 要 拍 同学 的 脸 — 只 拍 手\n• 测 试 时 让 别 的 同学 演示 — 看 AI 能 不 能 通 用")


# ============================================================
# TM-7 · OPTION B — 文具 识别 模型
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "✏️ Option B · 文具 识别 模型 · School Supplies AI", ORANGE)

# LEFT: 3 pair options
panel(s, 0.40, 0.95, 4.55, 4.10, ORANGE, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, ORANGE, "🎯 选 一 对 文具  Pick a Pair", sz=12)
pair_opts_b = [
    ("✏️ vs 🖍️", "铅笔 vs 马克 笔"),
    ("🧽 vs 📎", "橡皮 vs 胶 棒"),
    ("📕 vs 📓", "书 vs 本子"),
]
for i, (visual, cn) in enumerate(pair_opts_b):
    y = 1.55 + i*1.05
    pc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(4.25), Inches(0.95))
    pc.fill.solid(); pc.fill.fore_color.rgb = WARM; pc.line.fill.background()
    tb(s, 0.65, y+0.10, 4.05, 0.45, visual, sz=24, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
    tb(s, 0.65, y+0.55, 4.05, 0.32, cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# RIGHT TOP: data
panel(s, 5.15, 0.95, 4.45, 2.10, CYBER, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 0.95, 4.45, CYBER, "📷 收集 数据 怎么 拍?", sz=11)
data_tips_b = [
    "• 每 类 至少 20 张",
    "• 不 同 颜色",
    "• 不同 角度",
    "• 放 在 不同 位置",
]
for i, tip in enumerate(data_tips_b):
    tb(s, 5.30, 1.50 + i*0.35, 4.15, 0.30, tip, sz=10, b=True, c=DARK)

# RIGHT BOT: test + reflection
panel(s, 5.15, 3.20, 4.45, 1.85, PINK, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 3.20, 4.45, PINK, "🧪 测 试 + 反 思", sz=11)
tb(s, 5.30, 3.78, 4.15, 0.28, "🤔 测试 问 题:", sz=10, b=True, c=PINK)
tb(s, 5.30, 4.05, 4.15, 0.25, "• 换 不同 颜色 — 还 认 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.30, 4.15, 0.25, "• 都 很 长 — 会 搞 错 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.62, 4.15, 0.25, "💡 AI 看 什么? 形状? 长度? 颜色?", sz=10, b=True, c=PINK, a=PP_ALIGN.LEFT)

tb(s, 0.4, 5.15, 9.2, 0.25, "📦 学生 自带 — 投 入 度 高!",
   sz=10, b=True, c=ORANGE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "如 选 B:\n• 让 学生 提 前 知 道 — 带 自 己 的 文具\n• 「同 类 不同 颜色」 是 关 键: 让 AI 学 形状 而 不 是 颜色\n• 反 思: 如 果 都 是 黄 色, AI 可能 学 错 (只 认 颜色)")


# ============================================================
# TM-8 · OPTION C — 颜色 物品 分 类
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🌈 Option C · 颜色 物品 分 类 · Color AI", PINK)

# LEFT: 2 color pair options
panel(s, 0.40, 0.95, 4.55, 2.50, PINK, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, PINK, "🎯 选 一 对 颜色  Pick a Pair", sz=12)
pair_opts_c = [
    ("🔴 vs 🔵", "红 色 物 品 vs 蓝色 物品"),
    ("🟡 vs 🟢", "黄 色 物品 vs 绿色 物品"),
]
for i, (visual, cn) in enumerate(pair_opts_c):
    y = 1.55 + i*0.85
    pc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(4.25), Inches(0.75))
    pc.fill.solid(); pc.fill.fore_color.rgb = WARM; pc.line.fill.background()
    tb(s, 0.65, y+0.05, 4.05, 0.40, visual, sz=22, b=True, c=PINK, a=PP_ALIGN.CENTER)
    tb(s, 0.65, y+0.45, 4.05, 0.28, cn, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT BOT: IMPORTANT WARNING
wa = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.55), Inches(4.55), Inches(1.50))
wa.fill.solid(); wa.fill.fore_color.rgb = RED
wa.line.color.rgb = STAR; wa.line.width = Pt(3)
tb(s, 0.55, 3.65, 4.30, 0.30, "⚠️ 重要 警告!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.95, 4.30, 0.32, "不 要 只 用 1 支 红 铅笔 + 1 支 蓝 铅笔!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.30, 4.30, 0.30, "✅ 用 多 种 红 物品 + 多 种 蓝 物品",
   sz=11, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.65, 4.30, 0.28, "Otherwise AI just memorizes ONE object",
   sz=8, c=WARM, a=PP_ALIGN.LEFT)

# RIGHT TOP: data collection
panel(s, 5.15, 0.95, 4.45, 1.95, CYBER, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 0.95, 4.45, CYBER, "📷 收集 数据 怎么 拍?", sz=11)
data_tips_c = [
    "• 每 类 至少 20 张",
    "• 多 种 不同 物 品",
    "• 同 色 不 同 形状",
    "• 光 线 要 亮",
]
for i, tip in enumerate(data_tips_c):
    tb(s, 5.30, 1.45 + i*0.32, 4.15, 0.28, tip, sz=10, b=True, c=DARK)

# RIGHT BOT: test
panel(s, 5.15, 3.05, 4.45, 2.00, GREEN, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 3.05, 4.45, GREEN, "🧪 测 试 + 反 思", sz=11)
tb(s, 5.30, 3.62, 4.15, 0.25, "• 新 的 红 物品 — 还 认 吗?", sz=9, c=DARK)
tb(s, 5.30, 3.88, 4.15, 0.25, "• AI 学 的 是 颜色, 还是 物品?", sz=9, c=DARK)
tb(s, 5.30, 4.15, 4.15, 0.25, "• 光 线 暗 — 还 准 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.50, 4.15, 0.30, "💡 AI 可能 看: 颜色? 背景? 形状?",
   sz=10, b=True, c=GREEN, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "如 选 C:\n• 关 键 教 学 点: 数据 多 样性 — 不 能 只 用 1 个 物品\n• 否 则 AI 就 「死 记」 那 一 个 物品\n• 测 试 时 用 新 的 红 / 蓝 物品 — 看 AI 是 否 真 的 学 到 「颜色」\n• 联 系 Session 1 的 「白 色 = 兔子」 笑 话")


# ============================================================
# TM-9 · OPTION D — 形状 识别 模型
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔷 Option D · 形状 识别 模型 · Shape AI", PURPLE)

# LEFT: 2 shape pair options
panel(s, 0.40, 0.95, 4.55, 2.50, PURPLE, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, PURPLE, "🎯 选 一 对 形状  Pick a Pair", sz=12)
pair_opts_d = [
    ("⭕ vs 🔺", "圆 形 vs 三 角 形"),
    ("⬜ vs ▭", "正方 形 vs 长方 形"),
]
for i, (visual, cn) in enumerate(pair_opts_d):
    y = 1.55 + i*0.85
    pc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(4.25), Inches(0.75))
    pc.fill.solid(); pc.fill.fore_color.rgb = WARM; pc.line.fill.background()
    tb(s, 0.65, y+0.05, 4.05, 0.40, visual, sz=22, b=True, c=PURPLE, a=PP_ALIGN.CENTER)
    tb(s, 0.65, y+0.45, 4.05, 0.28, cn, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT BOT: prep tip
pt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.55), Inches(4.55), Inches(1.50))
pt.fill.solid(); pt.fill.fore_color.rgb = WARM
pt.line.color.rgb = PURPLE; pt.line.width = Pt(2.5)
tb(s, 0.55, 3.65, 4.30, 0.30, "🎨 准备 形状:",
   sz=12, b=True, c=PURPLE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.95, 4.30, 0.30, "• 画 / 剪 / 打 印 多 种 形状",
   sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.30, 4.30, 0.30, "• 不同 大 小 + 不同 颜色",
   sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.65, 4.30, 0.28, "• 旋 转 角度 也 要 试!",
   sz=10, b=True, c=PURPLE, a=PP_ALIGN.LEFT)

# RIGHT TOP: data collection
panel(s, 5.15, 0.95, 4.45, 1.95, CYBER, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 0.95, 4.45, CYBER, "📷 收集 数据 怎么 拍?", sz=11)
data_tips_d = [
    "• 每 类 至少 20 张",
    "• 不同 大 小",
    "• 不同 颜色",
    "• 旋 转 一 下 形状",
]
for i, tip in enumerate(data_tips_d):
    tb(s, 5.30, 1.45 + i*0.32, 4.15, 0.28, tip, sz=10, b=True, c=DARK)

# RIGHT BOT: test
panel(s, 5.15, 3.05, 4.45, 2.00, PINK, fill=WHITE, lw=2.5)
panel_head(s, 5.15, 3.05, 4.45, PINK, "🧪 测 试 + 反 思", sz=11)
tb(s, 5.30, 3.62, 4.15, 0.25, "• 换 颜色 — 还 认 吗?", sz=9, c=DARK)
tb(s, 5.30, 3.88, 4.15, 0.25, "• 旋 转 — 还 认 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.15, 4.15, 0.25, "• 变 大 / 小 — 还 认 吗?", sz=9, c=DARK)
tb(s, 5.30, 4.50, 4.15, 0.30, "💡 AI 看: 边 数? 角? 长 宽 比?",
   sz=10, b=True, c=PINK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "如 选 D:\n• 难 度 比 较 高 — 需要 准 备 形状\n• 老师 提前 准 备 卡 纸 / 打 印 / 剪 刀\n• 学 生 也 可以 画\n• 关 键 测 试: 旋 转 / 颜色 / 大小 — AI 是 否 真 学 到 「形状」")


# ============================================================
# TM-10 · GOOD DATA vs BAD DATA
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎯 好 数据 长 什么 样? · Good Data vs Bad Data", DAY)

# 2 columns: GOOD vs BAD
# LEFT: GOOD
panel(s, 0.40, 0.95, 4.55, 4.10, GREEN, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, GREEN, "✅ 好 数据  Good Data", sz=13)
good_tips = [
    ("📷", "清 楚"),
    ("📚", "数 量 多"),
    ("📐", "角 度 不 同"),
    ("💡", "光 线 好"),
    ("🌈", "不 只 一 个 例子"),
]
for i, (em, txt) in enumerate(good_tips):
    y = 1.60 + i * 0.62
    tb(s, 0.55, y, 0.50, 0.50, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, 1.10, y+0.08, 3.75, 0.40, txt, sz=15, b=True, c=GREEN)

# RIGHT: BAD
panel(s, 5.05, 0.95, 4.55, 4.10, RED, fill=WHITE, lw=3)
panel_head(s, 5.05, 0.95, 4.55, RED, "❌ 坏 数据  Bad Data", sz=13)
bad_tips = [
    ("📉", "太 少"),
    ("🌫️", "太 模 糊"),
    ("🗑️", "背 景 太 乱"),
    ("🔁", "两 类 太 像"),
    ("📋", "每 张 都 一 模 一 样"),
]
for i, (em, txt) in enumerate(bad_tips):
    y = 1.60 + i * 0.62
    tb(s, 5.20, y, 0.50, 0.50, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, 5.75, y+0.08, 3.75, 0.40, txt, sz=15, b=True, c=RED)

# Bottom takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.15), Inches(9.20), Inches(0.32))
tk.fill.solid(); tk.fill.fore_color.rgb = INK; tk.line.fill.background()
tb(s, 0.55, 5.18, 9.0, 0.25, "💡 AI 学 到 什么 — 取 决于 你 给 它 什么 数据!",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 分钟 重要 提醒:\n• 收集 数据 前 先 看 这 一 张\n• 「AI 不 是 魔法」 — 学 到 什么 取 决于 数据\n• 老师 强 调: 这 是 ML 最 重要 的 原 则")


# ============================================================
# TM-11 · TESTING RULES — 测 试 你 的 AI
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🧪 测 试 你 的 AI · Test Your AI", DAY)

# LEFT: 4 testing methods
panel(s, 0.40, 0.95, 4.55, 4.10, DAY, fill=WHITE, lw=3)
panel_head(s, 0.40, 0.95, 4.55, DAY, "📋 4 种 测试 · 4 Tests", sz=12)
test_methods = [
    ("1.", "用 训练 时 见 过 的", "Same as training"),
    ("2.", "用 新 的 例子", "Use NEW examples"),
    ("3.", "换 一 个 同学 测", "Different student tests"),
    ("4.", "故意 给 难 一 点 的", "Give a tricky example"),
]
for i, (num, cn, en) in enumerate(test_methods):
    y = 1.60 + i*0.78
    tb(s, 0.55, y, 0.50, 0.45, num, sz=20, b=True, c=DAY, a=PP_ALIGN.LEFT)
    tb(s, 1.10, y+0.04, 3.75, 0.35, cn, sz=13, b=True, c=DARK)
    tb(s, 1.10, y+0.40, 3.75, 0.28, en, sz=9, c=GRAY)

# RIGHT: record table
panel(s, 5.15, 0.95, 4.45, 4.10, ORANGE, fill=WARM, lw=3)
panel_head(s, 5.15, 0.95, 4.45, ORANGE, "📝 记 一 记 · Record", sz=12)

# Score lines
records = [
    ("✅", "猜 对 了 几 次?", "Correct: _____"),
    ("❌", "猜 错 了 几 次?", "Wrong: _____"),
    ("🔍", "哪 里 最 容易 错?", "Where most wrong?"),
]
for i, (em, cn, en) in enumerate(records):
    y = 1.65 + i*1.05
    rb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), Inches(y), Inches(4.15), Inches(0.90))
    rb.fill.solid(); rb.fill.fore_color.rgb = WHITE
    rb.line.color.rgb = ORANGE; rb.line.width = Pt(1.5)
    tb(s, 5.40, y+0.08, 0.55, 0.45, em, sz=22, a=PP_ALIGN.LEFT)
    tb(s, 5.95, y+0.10, 3.40, 0.35, cn, sz=13, b=True, c=ORANGE)
    tb(s, 5.95, y+0.48, 3.40, 0.30, en, sz=10, c=GRAY)

# Bottom prompt
tb(s, 0.4, 5.15, 9.2, 0.25, "📓 把 结果 写 下来 — 等 一 下 要 汇 报!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 测试 阶 段:\n• 4 种 测试 — 全 部 都 要 做\n• 重 点 是 「用 新 的 例子」 + 「换 同学 测」\n• 老师 提 醒: 记 录 「错 在 哪」 — 这 是 改 进 的 依 据")


# ============================================================
# TM-12 · IMPROVE YOUR MODEL — 改 进
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🔧 如 果 AI 猜 错 了 怎么 办? · How to Improve?", PINK)

tb(s, 0.4, 0.85, 9.2, 0.30, "5 个 改 进 办法 — 试 一 试!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 5 improvement methods
improvements = [
    ("➕", "加 更 多 照片", "Add more photos", CYBER),
    ("📷", "把 照片 拍 清 楚", "Make photos clearer", ORANGE),
    ("📐", "用 不同 角 度", "Try different angles", PURPLE),
    ("🗑️", "去 掉 容易 搞 错 的", "Remove confusing photos", PINK),
    ("✨", "找 更 好 的 例子", "Use better examples", GREEN),
]
iw = 1.78; igap = 0.12
itotal = 5*iw + 4*igap; istart = (10 - itotal)/2
for i, (em, cn, en, cl) in enumerate(improvements):
    x = istart + i*(iw + igap)
    panel(s, x, 1.35, iw, 2.65, cl, fill=WHITE, lw=2.5)
    tb(s, x, 1.50, iw, 0.85, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 2.45, iw, 0.40, cn, sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.88, iw-0.10, 0.28, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Number badge
    nbg = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+iw/2-0.22), Inches(3.45), Inches(0.45), Inches(0.45))
    nbg.fill.solid(); nbg.fill.fore_color.rgb = cl; nbg.line.fill.background()
    tb(s, x+iw/2-0.22, 3.50, 0.45, 0.35, str(i+1), sz=16, b=True, c=WHITE, a=PP_ALIGN.CENTER)

# Big teacher message
tm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.10))
tm.fill.solid(); tm.fill.fore_color.rgb = DAY
tm.line.color.rgb = STAR; tm.line.width = Pt(3)
tb(s, 0.55, 4.45, 9.0, 0.40, "💡 真 正 的 AI 工程师 不 是 一 次 就 成 功!",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.90, 9.0, 0.30, "他 们 测 试 → 发 现 问 题 → 再 改 进 → 再 测 试",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22, "Real AI engineers test, find problems, improve, retest!",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 改 进 阶 段:\n• 让 学 生 选 1-2 种 改 进 办法\n• 重新 训 练 → 再 测 试 → 看 有 没 有 进 步\n• 老师 鼓 励: 「错 了 不 要 怕 — 这 是 学 习 的 一 部分」")


# ============================================================
# TM-13 · GROUP REPORT — 小 组 汇 报 (sentence frames)
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🎤 小 组 汇 报 · Group Report", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "用 句 型 把 你 的 项 目 说 出 来!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Sentence frames in a big dark box
fp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.25), Inches(9.20), Inches(3.05))
fp.fill.solid(); fp.fill.fore_color.rgb = INK
fp.line.color.rgb = STAR; fp.line.width = Pt(3)
tb(s, 0.55, 1.35, 9.0, 0.30, "💬 我 们 的 模型:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)

# 6 sentence frame lines
report_frames = [
    "1. 我 们 的 模型 是: ______",
    "2. 它 能 分 辨: ______ 和 ______",
    "3. 我 们 每 类 拍 了 ______ 张 图片",
    "4. AI 猜 对 了 ______ 次",
    "5. AI 最 容易 搞 错 的 是 ______",
    "6. 我们 发 现: AI 需 要 ______",
]
for i, frame in enumerate(report_frames):
    y = 1.70 + i*0.42
    tb(s, 0.65, y, 8.80, 0.38, frame, sz=14, b=True, c=STAR, a=PP_ALIGN.LEFT)

# Bonus sentence frames
bf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.45), Inches(9.20), Inches(1.00))
bf.fill.solid(); bf.fill.fore_color.rgb = WARM
bf.line.color.rgb = DAY; bf.line.width = Pt(2)
tb(s, 0.55, 4.52, 9.0, 0.28, "✨ 加 分 句 型 · Bonus Frames:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.80, 9.0, 0.30, "• AI 看 到 了 ______ 特征",
   sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.12, 9.0, 0.30, "• AI 犯 错 是 因为 ______ ; 我们 改 进 的 方法 是 ______",
   sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "5-7 分钟 汇 报:\n• 每 组 2 分钟\n• 用 句 型 — 别 让 学 生 自由 发 挥 (容易 离 题)\n• 鼓 励 用 中文 — 数 字 / 「图片」 / 「特征」 都 是 关 键 词\n• 老师 在 板 上 记 录 各 组 数 据 — 比 较 哪 组 数据 多")


# ============================================================
# TM-14 · GALLERY TEST — AI 模型 挑 战 赛
# ============================================================
s = ns(prs); bg(s, CREAM, prs); hb(s, "🏆 AI 模型 挑 战 赛 · Gallery Test", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "请 别 的 组 来 测 试 你 的 模型!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 rules in 2x2 grid
gallery_rules = [
    ("🤝", "尊 重 别 人 的 模型", "Be respectful", CYBER),
    ("✅", "试 一 个 正 常 测试", "Try one normal test", GREEN),
    ("🎯", "试 一 个 难 一 点 的", "Try one tricky test", ORANGE),
    ("💡", "给 一 个 改 进 建 议", "Give one suggestion", PURPLE),
]
gw = 4.40; ggap_x = 0.20; ggap_y = 0.20; gh = 1.50
gstart_x = (10 - 2*gw - ggap_x)/2
for i, (em, cn, en, cl) in enumerate(gallery_rules):
    row = i // 2; col = i % 2
    x = gstart_x + col*(gw + ggap_x)
    y = 1.30 + row*(gh + ggap_y)
    panel(s, x, y, gw, gh, cl, fill=WHITE, lw=3)
    tb(s, x+0.15, y+0.20, 0.90, 1.00, em, sz=44, a=PP_ALIGN.LEFT)
    tb(s, x+1.20, y+0.25, gw-1.35, 0.45, cn, sz=15, b=True, c=cl)
    tb(s, x+1.20, y+0.75, gw-1.35, 0.32, en, sz=10, c=GRAY)
    # Number badge top-right
    nbg = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+gw-0.55), Inches(y+0.15), Inches(0.40), Inches(0.40))
    nbg.fill.solid(); nbg.fill.fore_color.rgb = cl; nbg.line.fill.background()
    tb(s, x+gw-0.55, 0.20+y, 0.40, 0.35, str(i+1), sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)

# Big finale message
fm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.50), Inches(9.20), Inches(0.95))
fm.fill.solid(); fm.fill.fore_color.rgb = DAY
fm.line.color.rgb = STAR; fm.line.width = Pt(3)
tb(s, 0.55, 4.62, 9.0, 0.45, "🤖 今天 你 们 都 是 AI 工程师 + AI 测试 员!",
   sz=17, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.30, "You're all AI engineers AND AI testers today!",
   sz=11, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "5 分钟 互 测:\n• 各 组 互 相 串 门 — A 组 测 B 组 的 模型\n• 注意 时 间 控制\n• 每 组 写 一 张 「改 进 建 议」 卡 给 对 方\n• 老师 收 集 卡 — 公 开 分 享 最 好 的 几 张")


# ============================================================
# 24 · SHARE + CLOSE
# ============================================================
s = share_close(prs, DAY,
    frames_cn=["「我们 训练 的 AI 会 认 ______」",
               "「数据 越 多, AI 就 ______」"],
    frames_en="Our AI can recognize ___  ·  More data makes AI ___",
    next_day_cn="Day 4 · 科技 改变 生活 — 以前 vs 现在!",
    next_day_en="Day 4 · How tech changes life — then vs now!",
    next_emoji="📱")
n += 1; pn(s, n)
notes(s, "10 分钟:\n• 每 组 1-2 分钟 分享\n• 句型 帮 学生 说 出 「数据 多 少」 + 「好 数据 vs 坏 数据」\n• 收 拾 教 具 准备 Day 4")


# ============================================================
# HIDDEN TEACHER NOTES — Emotion AI Teachable Machine guide
# Marked as hidden so it doesn't appear in slideshow mode.
# Single dense reference page for the teacher.
# ============================================================
s = ns(prs); bg(s, CREAM, prs)
# Header bar (PINK to signal "teacher only")
hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.55))
hd.fill.solid(); hd.fill.fore_color.rgb = PINK; hd.line.fill.background()
tb(s, 0.4, 0.20, 7.0, 0.45, "🧑‍🏫 TEACHER NOTES · Emotion AI Lab Guide",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.LEFT)
# HIDDEN badge
hb_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.60), Inches(0.23), Inches(2.00), Inches(0.40))
hb_box.fill.solid(); hb_box.fill.fore_color.rgb = STAR; hb_box.line.fill.background()
tb(s, 7.60, 0.28, 2.00, 0.32, "🙈 HIDDEN · 不 给 学 生 看", sz=10, b=True, c=INK, a=PP_ALIGN.CENTER)

# ROW 1: Goal + Before-class checklist (side-by-side)
# LEFT: Goal box
gb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(0.85), Inches(4.65), Inches(1.00))
gb.fill.solid(); gb.fill.fore_color.rgb = WARM; gb.line.color.rgb = DAY; gb.line.width = Pt(1.5)
tb(s, 0.40, 0.90, 4.45, 0.25, "🎯 Lesson Goal — Students should understand:",
   sz=10, b=True, c=DAY, a=PP_ALIGN.LEFT)
goal_lines = [
    "✅ AI 从 例子 学 习  · learns from examples",
    "✅ 更 多 数据 = 更 好 · more data helps",
    "✅ AI 会 犯错  · AI makes mistakes",
    "✅ 人 训练 + 改 进 AI  · humans improve AI",
]
for i, line in enumerate(goal_lines):
    tb(s, 0.45, 1.18 + i*0.17, 4.40, 0.18, line, sz=8, b=True, c=DARK, a=PP_ALIGN.LEFT)

# RIGHT: Before-class checklist
cb_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.05), Inches(0.85), Inches(4.65), Inches(1.00))
cb_box.fill.solid(); cb_box.fill.fore_color.rgb = WARM; cb_box.line.color.rgb = ORANGE; cb_box.line.width = Pt(1.5)
tb(s, 5.15, 0.90, 4.45, 0.25, "💻 BEFORE CLASS CHECKLIST:",
   sz=10, b=True, c=ORANGE, a=PP_ALIGN.LEFT)
check_lines = [
    "☐ Laptop + Chrome  ☐ 投影 仪 (Projector)",
    "☐ 摄像 头 工作 (Webcam works)  ☐ 网络 (Internet)",
    "☐ 提 前 打 开 + 测 试 一 次",
]
for i, line in enumerate(check_lines):
    tb(s, 5.20, 1.18 + i*0.20, 4.40, 0.20, line, sz=8, b=True, c=DARK, a=PP_ALIGN.LEFT)
# URL — hyperlink
url_tb = s.shapes.add_textbox(Inches(5.20), Inches(1.62), Inches(4.40), Inches(0.20))
url_tf = url_tb.text_frame; url_p = url_tf.paragraphs[0]
url_run = url_p.add_run()
url_run.text = "🔗 teachablemachine.withgoogle.com"
url_run.font.size = Pt(9); url_run.font.bold = True
url_run.font.color.rgb = CYBER; url_run.font.name = 'KaiTi'
url_run.hyperlink.address = "https://teachablemachine.withgoogle.com/"

# ROW 2: 9 STEPS in 2 columns (3 wide, 3 tall + 1 extra)
# Step header
tb(s, 0.30, 1.95, 9.40, 0.25, "📋 9 STEP WALKTHROUGH (~30-35 min total):",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)

steps_data = [
    ("1", "📂 Open Tool", "TM网址 → Get Started → Image Project → Standard", CYBER),
    ("2", "🏷️ Create Classes", "Class 1→😀开心, Class 2→😡生气, Class 3→😲惊讶, (Class 4→😢难过)", ORANGE),
    ("3", "📷 Round 1 (少 数据)", "每 类 仅 3-5 张 · 脸 清 楚, 表情 夸 张, 光 线 好", PURPLE),
    ("4", "🧠 Train Model", "点 Train Model → 等 1-2 分钟", DAY),
    ("5", "🧪 Test Round 1", "1-2 学 生 做 表情 → 问: AI 猜 对 吗? 为什么 错?", PINK),
    ("6", "📷 Round 2 (多 数据)", "每 类 20-30 张 · 不 同 学 生 / 角度 / 远 近 / 强度", GREEN),
    ("7", "🧠 Train Again", "再 点 Train → 「现在 AI 更 聪 明 吗?」", DAY),
    ("8", "🎯 Challenge Round", "😏 假 笑 / 🙃 怪 角度 / 🧢 戴 帽 / 😮 半 脸 / 🌙 暗 光", PINK),
    ("9", "🎓 Wrap-up", "复 习 6 步: 收集→标签→学习→找规律→测试→做判断", DAY),
]
# 2 columns: 5 left + 4 right
for i, (num, title, detail, cl) in enumerate(steps_data):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = 0.30 + col * 4.75
    y = 2.25 + row * 0.42
    # Step row
    rb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.65), Inches(0.38))
    rb.fill.solid(); rb.fill.fore_color.rgb = WHITE; rb.line.color.rgb = cl; rb.line.width = Pt(1)
    # Number circle
    nc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.05), Inches(y+0.07), Inches(0.25), Inches(0.25))
    nc.fill.solid(); nc.fill.fore_color.rgb = cl; nc.line.fill.background()
    tb(s, x+0.05, y+0.08, 0.25, 0.22, num, sz=10, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Title + detail
    tb(s, x+0.35, y+0.03, 4.25, 0.20, title, sz=9, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x+0.35, y+0.20, 4.25, 0.18, detail, sz=7, c=DARK, a=PP_ALIGN.LEFT)

# ROW 3: Troubleshooting + Timing
# LEFT: Troubleshooting
tr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(4.45), Inches(5.65), Inches(1.05))
tr.fill.solid(); tr.fill.fore_color.rgb = WARM; tr.line.color.rgb = RED; tr.line.width = Pt(1.5)
tb(s, 0.40, 4.50, 5.45, 0.22, "🔧 TROUBLESHOOTING:",
   sz=10, b=True, c=RED, a=PP_ALIGN.LEFT)
tr_lines = [
    "📷 摄像 头 不 工作 → 刷 新 / 允 许 摄像 头 权 限 / 重 启 Chrome",
    "🎯 准 确 度 低 → 加 更 多 图 / 改 善 光 线 / 表情 更 夸 张",
    "🐌 网 络 慢 → 减 少 到 3 个 类 别",
]
for i, line in enumerate(tr_lines):
    tb(s, 0.45, 4.75 + i*0.22, 5.40, 0.22, line, sz=8, b=True, c=DARK, a=PP_ALIGN.LEFT)

# RIGHT: Timing guide
tg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(4.45), Inches(3.65), Inches(1.05))
tg.fill.solid(); tg.fill.fore_color.rgb = INK; tg.line.color.rgb = STAR; tg.line.width = Pt(1.5)
tb(s, 6.15, 4.50, 3.45, 0.22, "⏱️ TIMING GUIDE:",
   sz=10, b=True, c=STAR, a=PP_ALIGN.LEFT)
tm_lines = [
    "Setup intro: 5 min   ·   R1 train: 5-8 min",
    "R1 test: 5 min       ·   R2 train: 8-10 min",
    "Challenge: 5 min     ·   Discussion: 5 min",
    "▶ TOTAL: ~30-35 min",
]
for i, line in enumerate(tm_lines):
    color = STAR if i == 3 else WHITE
    size = 9 if i == 3 else 8
    bold = True
    tb(s, 6.15, 4.74 + i*0.18, 3.50, 0.18, line, sz=size, b=bold, c=color, a=PP_ALIGN.LEFT)

# Mark this slide as HIDDEN (won't show in slideshow / slide sorter shows it grayed out)
s.element.set('show', '0')

# No page number on teacher slide (or could add but skip)
# n += 1; pn(s, n)  # intentionally omitted — hidden slide
notes(s, "TEACHER REFERENCE SLIDE — HIDDEN:\n• 这 张 不 给 学 生 看 — 仅 老师 使 用\n• 在 PowerPoint 中 此 slide 标 记 为 「隐藏」 (Hide Slide)\n• 播 放 时 跳 过 — 但 在 编 辑 视 图 可 见\n• 教 师 可 打 印 这 张 当 教 案\n\nKey teacher scripts (要 大 声 说):\n• Step 1: 「今天 我 们 来 训练 一 个 真 正 的 AI!」\n• Step 2: 「这 些 就 是 AI 要 学 习 的 答 案」\n• Step 3: 「我们 先 少 教 一 点, 看看 AI 会 不 会 聪 明」\n• Step 5: 「AI 不 是 魔法, 它 需 要 学 习」\n• Step 6: 「真 正 的 AI 要 看 很 多 很 多 例子」\n• Step 7: 「现在 AI 会 不 会 更 聪 明?」\n• Step 9: 「这 就 是 机器 学 习」")


out = os.path.join(os.path.dirname(__file__), "day3_ml.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
