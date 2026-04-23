#!/usr/bin/env python3
"""
Day 5 International Food Culture PPT — Same structure as Day 1-4 v2.
Theme: 国际美食文化 International Food Culture
Session 1 (11:00-11:45): 做世界美食 - Pizza / 墨西哥玉米饼 Tacos / 印度馕 Naan
Session 2 (2:00-2:45): 复习 D1-D4 + Bamboozle 分组比赛
Session 3 (3:00-4:30): 制作世界美食 - Sushi / 花馍
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# Colors (consistent with Day 1-4)
TEAL = RGBColor(0x00,0x69,0x5C)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BLACK = RGBColor(0x33,0x33,0x33)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
CREAM = RGBColor(0xFF,0xFA,0xF0)
WARM = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xE8,0xE8,0xE8)
BLUE = RGBColor(0x19,0x76,0xD2)
GREEN = RGBColor(0x38,0x8E,0x3C)
CORAL = RGBColor(0xE8,0x5D,0x04)  # Day 5 accent — food theme
SAFFRON = RGBColor(0xF4,0x8C,0x06)
# Country colors for foods
ITALY = RGBColor(0x00,0x89,0x4F)
MEXICO = RGBColor(0x00,0x6D,0x47)
INDIA = RGBColor(0xE6,0x78,0x12)
JAPAN = RGBColor(0xBC,0x00,0x2D)
CHINA = RGBColor(0xDE,0x29,0x10)
ASIA = RGBColor(0xC0,0x39,0x2B)
AFRICA = RGBColor(0xA1,0x66,0x2F)
EUROPE = RGBColor(0x15,0x3E,0x90)
AMERICAS = RGBColor(0x3C,0x3B,0x6E)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=BLACK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=BLACK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=TEAL,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color);tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER);tb(s,1,2.8,8,0.8,sub,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);return s
def vs(title,bgc):
    s=ns();bg(s,bgc);tb(s,1,0.8,8,0.8,"🎬 看视频  Watch Video",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,1.8,8,0.5,title,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
    ib(s,1.5,2.5,7,2.5,"📷 插入视频截图或粘贴视频链接");tb(s,1,5.1,8,0.3,"🔗 视频链接: ____________________",sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    return s

def recipe_slide(fe,name_cn,name_en,country,color,ingredients,steps,img_lb):
    """A recipe page: country color header, two-column ingredients + steps, image placeholder."""
    s=ns();bg(s,CREAM)
    tb(s,0.3,0.15,9.4,0.6,f"{fe} {name_cn} {name_en} — {country}",sz=24,b=True,c=color,a=PP_ALIGN.CENTER)
    # Ingredients column
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.5),Inches(0.4))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.4,0.98,4.3,0.35,"🥣 食材 Ingredients",sz=15,b=True,c=WHITE)
    tf=tb(s,0.4,1.45,4.4,3.0,ingredients[0],sz=14,c=DARK)
    for ing in ingredients[1:]:
        ap(tf,ing,sz=14,c=DARK)
    # Steps column
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(0.95),Inches(4.7),Inches(0.4))
    sh2.fill.solid();sh2.fill.fore_color.rgb=color;sh2.line.fill.background()
    tb(s,5.1,0.98,4.5,0.35,"👩‍🍳 步骤 Steps",sz=15,b=True,c=WHITE)
    tf2=tb(s,5.1,1.45,4.6,3.0,steps[0],sz=14,c=DARK)
    for st in steps[1:]:
        ap(tf2,st,sz=14,c=DARK)
    # Image at bottom
    ib(s,0.3,4.55,9.4,0.6,img_lb)
    return s

n=0

# ============================================================
# 1 COVER — Boarding pass for Day 5
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.8,"Global Explorer Camp",sz=36,b=True,c=TEAL,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.5,"环球探索沉浸式夏令营",sz=20,c=TEAL,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2),Inches(1.6),Inches(6),Inches(3.2));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CORAL;sh.line.width=Pt(3)
tf=tb(s,2.3,1.8,5.4,2.8,"BOARDING PASS  登机牌",sz=16,b=True,c=GRAY,a=PP_ALIGN.CENTER)
ap(tf,"",sz=8)
ap(tf,"Flight 航班: GR EDU-005",sz=18,c=DARK)
ap(tf,"Destination 目的地:  世界美食 WORLD CUISINE",sz=20,b=True,c=CORAL)
ap(tf,"Date 日期: June 12, 2025",sz=16,c=DARK)
ap(tf,"Gate 登机口: 谷雨厨房 GR EDU Kitchen",sz=16,c=DARK)
ap(tf,"",sz=8)
ap(tf,"Passenger 旅客: 谷雨小厨师们",sz=18,c=DARK)
ap(tf,"",sz=6)
ap(tf,"Ready to cook?  准备好做饭了吗？ 🍳",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","做美食 (3选1)：Pizza / Taco / Naan",CORAL),
    ("Session 2  下午","2:00-2:45","复习 D1-D4 + Bamboozle 分组比赛",BLUE),
    ("Session 3  下午","3:00-4:30","做美食：Sushi / 花馍 + 美食展览",GREEN)]):
    y=0.9+i*1.5;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2));sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE);tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=RGBColor(0xFF,0xF3,0xE0))
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=15,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"🍴 内容目标  Content:",sz=20,b=True,c=CORAL)
tf=tb(s,0.7,1.4,9,1.2,"1. 认识世界各国代表美食及其文化意义",sz=16,c=DARK)
ap(tf,"2. 学习厨房安全与卫生规则",sz=16,c=DARK)
ap(tf,"3. 动手制作5种国际美食，体验不同文化",sz=16,c=DARK)
tb(s,0.5,2.8,9,0.5,"🗣️ 语言目标  Language:",sz=20,b=True,c=BLUE)
tb(s,0.7,3.3,9,0.9,"📖 复习 D1-D4 重点生字词 + 国家名称 + 文化知识",sz=16,b=True,c=DARK)
tb(s,0.5,4.3,9,0.5,"🎉 实践目标  Practice: 完成 Bamboozle 比赛 + 美食展览",sz=16,c=GREEN)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","做世界美食 (三选一)   Pick 1 of 3\nPizza 🍕   Taco 🌮   Naan 🫓",CORAL,"🍳");n+=1

# ============================================================
# 5 WORLD FOOD MAP intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 世界美食地图  World Food Map",CORAL)
tb(s,0.4,0.9,4.5,0.5,"每个国家都有自己的代表美食！",sz=18,b=True,c=CORAL)
tb(s,0.4,1.4,4.5,0.4,"Every country has its iconic dish!",sz=14,c=GRAY)
tf=tb(s,0.4,2.0,4.5,3.0,"🍕 意大利 Italy — Pizza",sz=16,c=DARK)
ap(tf,"🌮 墨西哥 Mexico — Taco 玉米饼",sz=16,c=DARK)
ap(tf,"🫓 印度 India — Naan 馕",sz=16,c=DARK)
ap(tf,"🍣 日本 Japan — Sushi 寿司",sz=16,c=DARK)
ap(tf,"🌸 中国 China — 花馍 Flower Mantou",sz=16,c=DARK)
ib(s,5.2,0.9,4.5,4.2,"📷 世界美食地图")
pn(s,n)

# ============================================================
# 6 KITCHEN SAFETY
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧼 厨房安全与卫生  Kitchen Safety",CORAL)
rules=[
    ("🧼 洗手","做饭前先洗手 20 秒  Wash hands 20s"),
    ("👕 系围裙","穿围裙、扎头发  Apron + tie back hair"),
    ("🔪 小心刀","刀朝下，拿稳  Knife down, hold firm"),
    ("🔥 离火远","远离热源  Stay away from hot surfaces"),
    ("🧽 保持干净","用完擦桌子  Clean up after use"),
    ("🙋‍♀️ 问老师","不懂就举手  Ask teacher if unsure"),
]
for i,(t,d) in enumerate(rules):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=0.95+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=CORAL;sh.line.width=Pt(2)
    tb(s,x+0.15,y+0.1,2.8,0.5,t,sz=18,b=True,c=CORAL,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,y+0.7,2.8,1.0,d,sz=13,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 7 CHOICE OVERVIEW — Pick your station
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🍽️ 选一个站点  Pick Your Station (3选1)",CORAL)
stations=[
    ("🍕","Pizza","意大利 Italy","自己做饼皮\n加酱料+奶酪+配料\n烤 10 分钟",ITALY),
    ("🌮","Taco 玉米饼","墨西哥 Mexico","加肉、蔬菜\n+牛油果酱+奶酪\n卷起来吃！",MEXICO),
    ("🫓","Naan 馕","印度 India","揉面团\n拍成圆形\n煎到金黄",INDIA),
]
for i,(em,nm,c,d,cl) in enumerate(stations):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(3.1),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.1,2.9,0.7,em,sz=48,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.95,2.9,0.45,nm,sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.45,2.9,0.35,c,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.15,2.9,2.85,1.8,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=13,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8 PIZZA intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🍕 Choice 1: Pizza 披萨 — 意大利 Italy",ITALY)
tb(s,0.4,0.9,4.5,0.5,"🇮🇹 你知道吗？  Did you know?",sz=18,b=True,c=ITALY)
tf=tb(s,0.4,1.5,4.5,3.5,"🍅 披萨起源于意大利那不勒斯",sz=14,c=DARK)
ap(tf,"Pizza was born in Naples, Italy!",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🧀 最经典的是玛格丽塔披萨",sz=14,c=DARK)
ap(tf,"Margherita: 番茄红 + 奶酪白 + 罗勒绿 = 意大利国旗！",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🔥 传统用木柴烤炉烤制",sz=14,c=DARK)
ap(tf,"Cooked in wood-fired oven at 900°F!",sz=12,c=GRAY)
ib(s,5.2,0.9,4.5,4.2,"📷 Pizza / Naples")
pn(s,n)

# ============================================================
# 9 PIZZA recipe
# ============================================================
recipe_slide("🍕","Pizza","披萨","意大利 Italy",ITALY,
    ["• 披萨饼皮 Pizza dough (已准备)","• 番茄酱 Tomato sauce","• 马苏里拉奶酪 Mozzarella","• 配料 Toppings: 火腿、蔬菜、橄榄","• 罗勒叶 Basil (可选)"],
    ["1️⃣ 饼皮上刷番茄酱","2️⃣ 撒上奶酪","3️⃣ 加你喜欢的配料","4️⃣ 老师帮忙放入烤箱","5️⃣ 425°F 烤 10-12 分钟","6️⃣ 取出晾 1 分钟，切块！"],
    "📷 Pizza 步骤图")
n+=1;pn(s,n)

# ============================================================
# 10 TACO intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌮 Choice 2: Taco 墨西哥玉米饼 — Mexico",MEXICO)
tb(s,0.4,0.9,4.5,0.5,"🇲🇽 你知道吗？  Did you know?",sz=18,b=True,c=MEXICO)
tf=tb(s,0.4,1.5,4.5,3.5,"🌽 玉米饼 (Tortilla) 是墨西哥主食",sz=14,c=DARK)
ap(tf,"Corn tortillas — Mexico's daily bread",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"📜 玉米饼有 9000 年历史！",sz=14,c=DARK)
ap(tf,"Over 9,000 years old — from Aztec times!",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🎉 Taco Tuesday 是美国超流行的传统",sz=14,c=DARK)
ap(tf,"Every Tuesday = Taco day in the USA!",sz=12,c=GRAY)
ib(s,5.2,0.9,4.5,4.2,"📷 Taco / Mexico")
pn(s,n)

# ============================================================
# 11 TACO recipe
# ============================================================
recipe_slide("🌮","Taco","墨西哥玉米饼","Mexico",MEXICO,
    ["• 玉米饼 Corn tortilla (小)","• 鸡肉丝或牛肉末 Chicken/beef","• 生菜、番茄、洋葱 Veggies","• 奶酪丝 Shredded cheese","• 牛油果酱 Guacamole","• 莎莎酱 Salsa"],
    ["1️⃣ 在平底锅上热玉米饼 30 秒","2️⃣ 放在盘子上","3️⃣ 加肉 (先放热的)","4️⃣ 加蔬菜和奶酪","5️⃣ 浇上牛油果酱和莎莎酱","6️⃣ 对折，享受！¡Buen provecho!"],
    "📷 Taco 步骤图")
n+=1;pn(s,n)

# ============================================================
# 12 NAAN intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🫓 Choice 3: Naan 印度馕 — India",INDIA)
tb(s,0.4,0.9,4.5,0.5,"🇮🇳 你知道吗？  Did you know?",sz=18,b=True,c=INDIA)
tf=tb(s,0.4,1.5,4.5,3.5,"🌾 馕是印度最流行的饼",sz=14,c=DARK)
ap(tf,"Naan — the most popular Indian flatbread",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🔥 传统在 Tandoor (泥炉) 里烤",sz=14,c=DARK)
ap(tf,"Baked in a clay oven — Tandoor!",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🥘 通常配咖喱 Curry 吃",sz=14,c=DARK)
ap(tf,"Perfect with curry — tear and dip!",sz=12,c=GRAY)
ib(s,5.2,0.9,4.5,4.2,"📷 Naan / Tandoor")
pn(s,n)

# ============================================================
# 13 NAAN recipe
# ============================================================
recipe_slide("🫓","Naan","印度馕","India",INDIA,
    ["• 面团 Pre-made dough (小份)","• 面粉 Flour (撒台面用)","• 黄油 Melted butter","• 大蒜末 Garlic (可选)","• 香菜 Cilantro (可选)","• 盐 Salt 一点点"],
    ["1️⃣ 把面团压扁成椭圆形","2️⃣ 撒面粉防粘","3️⃣ 平底锅不放油，中火","4️⃣ 煎 1-2 分钟，起泡翻面","5️⃣ 再煎 1 分钟","6️⃣ 刷黄油+蒜末，趁热吃！"],
    "📷 Naan 步骤图")
n+=1;pn(s,n)

# ============================================================
# 14 Taste & Share
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"😋 分享品尝  Share & Taste",CORAL)
tb(s,0.4,0.9,9,0.5,"做完自己的，和同学交换尝一尝！",sz=18,b=True,c=CORAL,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Finish your dish, then swap with friends — try all three!",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
prompts=[
    ("🍕","Pizza","“好吃！Delizioso!” (意大利语)",ITALY),
    ("🌮","Taco","“好吃！¡Delicioso!” (西班牙语)",MEXICO),
    ("🫓","Naan","“好吃！स्वादिष्ट!” (印度语)",INDIA),
]
for i,(em,nm,say,cl) in enumerate(prompts):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(3.1),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,x+0.1,2.1,2.9,0.8,em,sz=42,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.0,2.9,0.4,nm,sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.5,2.9,1.4,say,sz=12,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 15 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","复习 D1-D4 重点知识 + 生字词\n🎮 Bamboozle 环球之旅分组比赛",BLUE,"📖");n+=1

# ============================================================
# 16-19 Review D1-D4 (one slide per continent)
# ============================================================
reviews=[
    # emoji, cn, en, color, capitals (country, capital, food, landmark)
    ("🌏","亚洲 Asia","Day 1",ASIA,[
        ("🇨🇳 中国","北京","饺子/花馍","长城/故宫"),
        ("🇯🇵 日本","东京","寿司/拉面","富士山"),
        ("🇰🇷 韩国","首尔","泡菜/烤肉","景福宫"),
    ]),
    ("🌍","非洲 Africa","Day 2",AFRICA,[
        ("🇪🇬 埃及","开罗","库沙利/面包","金字塔"),
        ("🇰🇪 肯尼亚","内罗毕","Ugali","马赛马拉"),
        ("🇿🇦 南非","比勒陀利亚","Braai 烤肉","好望角"),
    ]),
    ("🌍","欧洲 Europe","Day 3",EUROPE,[
        ("🇫🇷 法国","巴黎","可颂/马卡龙","埃菲尔铁塔"),
        ("🇮🇹 意大利","罗马","Pizza/意面","斗兽场"),
        ("🇬🇧 英国","伦敦","炸鱼薯条","大本钟"),
    ]),
    ("🌎","美洲 Americas","Day 4",AMERICAS,[
        ("🇺🇸 美国","华盛顿","汉堡/热狗","自由女神像"),
        ("🇨🇦 加拿大","渥太华","Poutine/枫糖","尼亚加拉瀑布"),
        ("🇲🇽 墨西哥","墨西哥城","Taco/牛油果酱","奇琴伊察"),
    ]),
]
for em,cont,day,color,rows in reviews:
    s=ns();n+=1;bg(s,CREAM);hb(s,f"{em} 复习 Review — {cont}  ({day})",BLUE)
    # table 4 cols x 4 rows (header + 3)
    ts=s.shapes.add_table(4,4,Inches(0.3),Inches(1.0),Inches(9.4),Inches(3.9));t=ts.table
    t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.0);t.columns[2].width=Inches(2.8);t.columns[3].width=Inches(2.8)
    hd=["国家 Country","首都 Capital","美食 Food","景点 Landmark"]
    for c,ct in enumerate(hd):
        cl=t.cell(0,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        rn=p.add_run();rn.text=ct;rn.font.name='KaiTi';rn.font.size=Pt(13);rn.font.bold=True;rn.font.color.rgb=WHITE
        cl.fill.solid();cl.fill.fore_color.rgb=color
    for r,row in enumerate(rows,start=1):
        for c,ct in enumerate(row):
            cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
            rn=p.add_run();rn.text=ct;rn.font.name='KaiTi';rn.font.size=Pt(12);rn.font.bold=(c==0);rn.font.color.rgb=DARK
            if r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF5)
            else:cl.fill.solid();cl.fill.fore_color.rgb=WARM
    tb(s,0.4,5.05,9,0.3,"👀 看图+听老师读 → 大声跟读！",sz=12,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# ============================================================
# 20 我会认 review flash
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👀 生字词复习  Vocabulary Review — 我会认",BLUE)
words_read=[
    ("亚洲","yà zhōu","Asia"),("中国","zhōng guó","China"),("日本","rì běn","Japan"),
    ("非洲","fēi zhōu","Africa"),("欧洲","ōu zhōu","Europe"),("意大利","yì dà lì","Italy"),
    ("美洲","měi zhōu","Americas"),("美国","měi guó","USA"),("巴西","bā xī","Brazil"),
    ("首都","shǒu dū","capital"),("国旗","guó qí","flag"),("美食","měi shí","food"),
]
for i,(w,py,en) in enumerate(words_read):
    col=i%4;row=i//4
    x=0.3+col*2.4;y=1.0+row*1.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.3),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=BLUE;sh.line.width=Pt(1.5)
    tb(s,x+0.1,y+0.08,2.1,0.55,w,sz=26,b=True,c=BLUE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.68,2.1,0.28,py,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.93,2.1,0.28,en,sz=11,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 21 我会写 review
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"✍️ 生字词复习  Vocabulary Review — 我会写",BLUE)
words_write=[("亚洲","yà zhōu","Asia"),("中国","zhōng guó","China"),("非洲","fēi zhōu","Africa"),("欧洲","ōu zhōu","Europe"),("美洲","měi zhōu","Americas")]
for i,(w,py,en) in enumerate(words_write):
    if i<3:x=0.3+i*3.15;y=1.0
    else:x=1.9+(i-3)*3.15;y=3.1
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BLUE;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.1,2.8,1.0,w,sz=52,b=True,c=BLUE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.15,2.8,0.35,py,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.5,2.8,0.35,en,sz=13,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,5.15,9,0.3,"📝 每个字写 3 遍",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 22 BAMBOOZLE intro
# ============================================================
s=ns();n+=1;bg(s,BLUE)
tb(s,1,0.8,8,0.8,"🎮 Bamboozle 大比拼",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"环球之旅分组比赛  Global Tour Team Competition",sz=20,c=WARM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(2.5),Inches(7),Inches(2.4))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tf=tb(s,1.7,2.7,6.6,0.4,"📋 规则  Rules:",sz=18,b=True,c=BLUE)
ap(tf,"1. 分成 4 队 (亚洲队、非洲队、欧洲队、美洲队)",sz=14,c=DARK)
ap(tf,"2. 每队轮流答题，答对得分",sz=14,c=DARK)
ap(tf,"3. 答错也没关系，下一队继续",sz=14,c=DARK)
ap(tf,"4. 题目涵盖 D1-D4 所有内容",sz=14,c=DARK)
ap(tf,"5. 最高分队获得「环球之旅冠军」奖杯 🏆",sz=14,b=True,c=CORAL)
pn(s,n)

# ============================================================
# 23 BAMBOOZLE sample Q
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 Bamboozle 样题  Sample Questions",BLUE)
qs=[
    ("Q1","中国的首都是哪里？  Capital of China?","北京 Beijing",ASIA),
    ("Q2","哪个国家有金字塔？  Which country has pyramids?","埃及 Egypt",AFRICA),
    ("Q3","埃菲尔铁塔在哪里？  Where is Eiffel Tower?","法国巴黎 Paris, France",EUROPE),
    ("Q4","Taco 是哪个国家的？  Taco is from?","墨西哥 Mexico",AMERICAS),
]
for i,(lbl,q,ans,cl) in enumerate(qs):
    col=i%2;row=i//2
    x=0.3+col*4.75;y=1.0+row*2.0
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.6),Inches(1.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,x+0.15,y+0.1,0.8,0.4,lbl,sz=16,b=True,c=cl)
    tb(s,x+0.15,y+0.55,4.3,0.5,q,sz=13,b=True,c=DARK)
    tb(s,x+0.15,y+1.15,4.3,0.4,f"✅ {ans}",sz=13,c=GREEN)
    tb(s,x+0.15,y+1.45,4.3,0.35,"— 点击揭晓  Click to reveal —",sz=10,c=GRAY)
pn(s,n)

# ============================================================
# 24 Scoreboard
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 分数板  Scoreboard",BLUE)
teams=[("🌏 亚洲队","Team Asia",ASIA),("🌍 非洲队","Team Africa",AFRICA),("🌍 欧洲队","Team Europe",EUROPE),("🌎 美洲队","Team Americas",AMERICAS)]
for i,(nm,en,cl) in enumerate(teams):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.2),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,x+0.1,1.1,2.0,0.5,nm,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.6,2.0,0.4,en,sz=12,c=WARM,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(2.3),Inches(1.8),Inches(1.5))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.fill.background()
    tb(s,x+0.2,2.5,1.8,1.0,"___",sz=48,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.0,2.0,0.4,"分 Points",sz=12,c=WARM,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.45,2.0,0.4,"🏆 冠军？",sz=12,b=True,c=WARM,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 25 SESSION 3 DIVIDER
# ============================================================
div("Session 3  下午","制作世界美食 (二选一)   Pick 1 of 2\n🍣 Sushi 寿司 (日本)   🌸 花馍 (中国)",GREEN,"🍱");n+=1

# ============================================================
# 26 SUSHI intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🍣 Choice 1: Sushi 寿司 — 日本 Japan",JAPAN)
tb(s,0.4,0.9,4.5,0.5,"🇯🇵 你知道吗？  Did you know?",sz=18,b=True,c=JAPAN)
tf=tb(s,0.4,1.5,4.5,3.5,"🐟 寿司源于古代保存鱼的方法",sz=14,c=DARK)
ap(tf,"Sushi started as a way to preserve fish!",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🍙 寿司 = 醋米饭 + 配料",sz=14,c=DARK)
ap(tf,"Sushi = vinegared rice + toppings",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🎎 常见种类：Nigiri / Maki / Temaki",sz=14,c=DARK)
ap(tf,"Today we make Maki rolls 卷寿司!",sz=12,c=GRAY)
ib(s,5.2,0.9,4.5,4.2,"📷 Sushi / Japan")
pn(s,n)

# ============================================================
# 27 SUSHI recipe
# ============================================================
recipe_slide("🍣","Sushi (Maki)","寿司卷","Japan",JAPAN,
    ["• 寿司米 Cooked sushi rice (已调味)","• 海苔 Nori seaweed 一片","• 黄瓜条 Cucumber","• 牛油果 Avocado","• 蟹肉棒 Imitation crab","• 竹帘 Bamboo mat"],
    ["1️⃣ 竹帘上铺海苔(光面朝下)","2️⃣ 均匀铺一层米饭 (留 1 英寸边)","3️⃣ 中间摆上黄瓜+牛油果+蟹肉","4️⃣ 用竹帘慢慢卷紧","5️⃣ 用湿刀切成 6-8 块","6️⃣ 配酱油 ごちそうさま!"],
    "📷 Sushi 步骤图")
n+=1;pn(s,n)

# ============================================================
# 28 HUAMO intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌸 Choice 2: 花馍 Flower Mantou — 中国 China",CHINA)
tb(s,0.4,0.9,4.5,0.5,"🇨🇳 你知道吗？  Did you know?",sz=18,b=True,c=CHINA)
tf=tb(s,0.4,1.5,4.5,3.5,"🌾 花馍是中国北方的传统面食",sz=14,c=DARK)
ap(tf,"Decorative steamed bread from North China",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🎊 节日、婚礼、生日才做",sz=14,c=DARK)
ap(tf,"Made for festivals, weddings, birthdays",sz=12,c=GRAY)
ap(tf,"",sz=6)
ap(tf,"🎨 山西花馍被列入国家非物质文化遗产",sz=14,c=DARK)
ap(tf,"Shanxi Huamo — National intangible heritage!",sz=12,c=GRAY)
ib(s,5.2,0.9,4.5,4.2,"📷 花馍 / 中国")
pn(s,n)

# ============================================================
# 29 HUAMO recipe
# ============================================================
recipe_slide("🌸","花馍","Flower Mantou","China",CHINA,
    ["• 面团 Pre-made dough","• 天然色素 Natural colors:","   红(甜菜)、黄(南瓜)、绿(菠菜)","• 红枣 Jujubes (装饰)","• 黑芝麻 Black sesame (眼睛)","• 蒸锅 Steamer"],
    ["1️⃣ 把面团分成小球","2️⃣ 加色素揉匀 (每色一份)","3️⃣ 捏成花、兔、鱼等形状","4️⃣ 用红枣/芝麻装饰","5️⃣ 蒸 15-20 分钟","6️⃣ 出锅晾凉，摆盘展示！"],
    "📷 花馍 步骤图")
n+=1;pn(s,n)

# ============================================================
# 30 Food Exhibition
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎉 美食展览  Food Exhibition",GREEN)
tb(s,0.4,0.9,9,0.5,"把你做的美食摆出来，和大家分享！",sz=18,b=True,c=GREEN,a=PP_ALIGN.CENTER)
tb(s,0.4,1.4,9,0.4,"Display all your dishes — taste, take photos, celebrate!",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
steps=[
    ("1️⃣","📸","拍照 Photo","拍一张你做的美食"),
    ("2️⃣","🏷️","贴标签 Label","写上名字+国家"),
    ("3️⃣","😋","尝一尝 Taste","尝尝同学的美食"),
    ("4️⃣","⭐","投票 Vote","最喜欢哪一道？"),
]
for i,(num,em,t,d) in enumerate(steps):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=GREEN;sh.line.width=Pt(2)
    tb(s,x+0.1,2.1,2.0,0.5,num,sz=22,b=True,c=GREEN,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.6,2.0,0.8,em,sz=40,c=GREEN,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.55,2.0,0.45,t,sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.05,2.0,0.8,d,sz=12,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 31 World Traveler Stamp — FINAL
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.4,8,0.8,"🪪 环球美食家证书  World Foodie Certificate",sz=28,b=True,c=CORAL,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.4),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CORAL;sh.line.width=Pt(5)
tf=tb(s,3.6,1.65,2.8,2.7,"WORLD\n FOODIE",sz=24,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=14,b=True,c=GREEN,a=PP_ALIGN.CENTER)
ap(tf,"6/12/2025",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
ap(tf,"🍕 🌮 🫓 🍣 🌸",sz=18,a=PP_ALIGN.CENTER)
tb(s,1,4.55,8,0.4,"恭喜你完成环球美食之旅！Congratulations, world chef! 🎉",sz=16,b=True,c=CORAL,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"5 天 · 15+ 个国家 · 5 种美食 · 无数回忆",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 32 Wrap-up / Thank you
# ============================================================
s=ns();n+=1;bg(s,TEAL)
tb(s,1,0.8,8,0.8,"✈️ 旅程结束  Journey Complete",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"Global Explorer Camp — Day 5",sz=18,c=WARM,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.6,7,2.5,"🌏 5 天走遍 5 大洲",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"5 days · 5 continents · 15+ countries",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🎨 手工 + 🍳 做饭 + 📖 生字 + 🎮 游戏",sz=16,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"Thank you, young explorers! 谢谢小探险家们！",sz=16,b=True,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"See you next summer! 明年夏天见！✨",sz=14,c=WARM,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/世界旅行world_trip_pbl/day5_food_culture.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
