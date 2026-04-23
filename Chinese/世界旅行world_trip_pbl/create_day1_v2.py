#!/usr/bin/env python3
"""
Day 1 Asia PPT v2 — Clean rebuild based on 3-session schedule.
Each slide: 2-3 items max, with image placeholder space.
Output: day1_asia_v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W = prs.slide_width
H = prs.slide_height

# ─── Colors ───
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
GRAY = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xBB, 0xBB, 0xBB)
CHINA_RED = RGBColor(0xDE, 0x29, 0x10)
JAPAN_RED = RGBColor(0xBC, 0x00, 0x2D)
INDIA_GREEN = RGBColor(0x13, 0x88, 0x08)
BG_CREAM = RGBColor(0xFF, 0xFA, 0xF0)
BG_WARM = RGBColor(0xFF, 0xF3, 0xE0)
IMG_BG = RGBColor(0xE8, 0xE8, 0xE8)

# ─── Helpers ───
def new_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def tb(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=None, font='KaiTi'):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align: p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tf

def add_p(tf, text, size=18, bold=False, color=BLACK, align=None):
    p = tf.add_paragraph()
    if align: p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = 'KaiTi'

def bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def img_box(slide, left, top, width, height, label="📷"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = IMG_BG
    shape.line.fill.background()
    tb(slide, left+0.1, top + height/2 - 0.2, width-0.2, 0.4, f"{label}", size=14, color=LGRAY, align=PP_ALIGN.CENTER)

def header_bar(slide, text, color=ORANGE, top=0.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(top), Inches(9.4), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tb(slide, 0.4, top+0.03, 9.2, 0.5, text, size=20, bold=True, color=WHITE)

def page_num(slide, num):
    tb(slide, 9.0, 5.25, 0.8, 0.3, str(num), size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def deco_corners(slide, emojis="🐼🌸🗻🐘"):
    """Add small decorative emojis in corners to make slides more playful"""
    tb(slide, 0.1, 5.0, 0.5, 0.4, emojis[0], size=16, color=LGRAY)
    tb(slide, 9.3, 5.0, 0.5, 0.4, emojis[-1], size=16, color=LGRAY)

def section_divider(title, subtitle, color, emoji=""):
    s = new_slide()
    bg(s, color)
    tb(s, 1, 1.5, 8, 1.2, f"{emoji} {title}", size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, 1, 2.8, 8, 0.8, subtitle, size=22, color=RGBColor(0xFF, 0xF3, 0xE0), align=PP_ALIGN.CENTER)
    return s

def country_info_slide(flag_emoji, cn, en, color, title, items, img_label="📷"):
    """Generic country info slide: title + 2-3 text items + image placeholder"""
    s = new_slide()
    bg(s, BG_CREAM)
    # Title
    tb(s, 0.3, 0.15, 9.4, 0.6, f"{flag_emoji} {cn} {en} — {title}", size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Text items (left side)
    y = 1.0
    for item_title, item_detail in items:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(y), Inches(4.8), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tb(s, 0.4, y+0.02, 4.6, 0.35, item_title, size=15, bold=True, color=WHITE)
        tb(s, 0.4, y+0.45, 4.7, 0.5, item_detail, size=16, color=DARK)
        y += 1.1
    # Image placeholder (right side)
    img_box(s, 5.5, 1.0, 4.2, y - 1.0 - 0.1, img_label)
    # Decorative Asian emoji (bottom-right corner)
    deco = {"中国": "🐼", "日本": "🌸", "印度": "🐅"}.get(cn, "🌏")
    tb(s, 8.8, 4.8, 1.0, 0.6, deco, size=28, color=LGRAY, align=PP_ALIGN.RIGHT)
    return s

n = 0  # slide counter

# ═══════════════════════════════════════════
# SLIDE 1: COVER / BOARDING PASS
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
tb(s, 1, 0.3, 8, 0.8, "Global Explorer Camp", size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, 1, 0.9, 8, 0.5, "环球探索沉浸式夏令营", size=20, color=ORANGE, align=PP_ALIGN.CENTER)
# Boarding pass box
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.6), Inches(6), Inches(3.2))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = ORANGE
shape.line.width = Pt(3)
tf = tb(s, 2.3, 1.8, 5.4, 2.8, "BOARDING PASS  登机牌", size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_p(tf, "", size=8)
add_p(tf, "Flight 航班: GR EDU-001", size=18, color=DARK)
add_p(tf, "Destination 目的地:  亚洲 ASIA", size=20, bold=True, color=ORANGE)
add_p(tf, "Date 日期: June 8, 2025", size=16, color=DARK)
add_p(tf, "Gate 登机口: 谷雨大厅 GR EDU Hall", size=16, color=DARK)
add_p(tf, "", size=8)
add_p(tf, "Passenger 旅客: 谷雨全体师生", size=18, color=DARK)
add_p(tf, "", size=6)
add_p(tf, "Fasten seatbelts!  系好安全带！ ✈️", size=14, color=GRAY, align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# SLIDE 2: 时间安排
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "⏰ 今日时间安排  Today's Schedule")
sessions = [
    ("Session 1  上午", "11:00-11:45", "了解亚洲 + 三个国家", ORANGE),
    ("Session 2  下午", "2:00-2:45", "复习总结 + 语言目标（认字写字）", RGBColor(0x19,0x76,0xD2)),
    ("Session 3  下午", "3:00-4:30", "写Booklet + 做Project", RGBColor(0x38,0x8E,0x3C)),
]
for i, (name, time, desc, clr) in enumerate(sessions):
    y = 0.9 + i * 1.5
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    tb(s, 0.7, y+0.1, 4, 0.4, name, size=20, bold=True, color=WHITE)
    tb(s, 0.7, y+0.5, 3, 0.4, time, size=16, color=RGBColor(0xFF,0xF3,0xE0))
    tb(s, 4.5, y+0.15, 4.8, 0.9, desc, size=18, color=WHITE)
page_num(s, n)

# ═══════════════════════════════════════════
# SLIDE 3: 教学目标
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🎯 教学目标  Learning Objectives")
# Content goals
tb(s, 0.5, 0.9, 9, 0.5, "📚 内容目标 Content Goals:", size=20, bold=True, color=ORANGE)
tf = tb(s, 0.7, 1.4, 8.5, 1.2, "1. 了解亚洲的地理位置和特点", size=16, color=DARK)
add_p(tf, "2. 了解三个国家：中国、日本、印度（国旗、首都、景点、文化）", size=16, color=DARK)
# Language goals
tb(s, 0.5, 2.8, 9, 0.5, "🗣️ 语言目标 Language Goals:", size=20, bold=True, color=RGBColor(0x19,0x76,0xD2))
tb(s, 0.7, 3.3, 4.2, 0.9, "👀 我会认：亚洲  中国  日本  印度  首都", size=17, bold=True, color=DARK)
tb(s, 5.2, 3.3, 4.2, 0.9, "✍️ 我会写：中国  日本  亚洲  首都", size=17, bold=True, color=DARK)
# Project goals
tb(s, 0.5, 4.3, 9, 0.5, "🎨 实践目标: 完成亚洲Booklet + 手工项目", size=16, color=RGBColor(0x38,0x8E,0x3C))
page_num(s, n)

# ═══════════════════════════════════════════
# SESSION 1: 上午 DIVIDER
# ═══════════════════════════════════════════
section_divider("Session 1  上午", "了解亚洲的地理位置和特点\n了解三个主要国家：中国、日本、印度", ORANGE, "🌏"); n+=1

# SLIDE: 亚洲视频
s = new_slide(); n+=1
bg(s, RGBColor(0x1A,0x23,0x7E))
tb(s, 1, 0.8, 8, 0.8, "🎬 看视频  Watch Video", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1, 1.8, 8, 0.5, "认识亚洲 About Asia", size=22, color=RGBColor(0xFF,0xF3,0xE0), align=PP_ALIGN.CENTER)
img_box(s, 1.5, 2.5, 7, 2.5, "📷 插入视频截图或粘贴视频链接\nInsert video screenshot or paste link here")
tb(s, 1, 5.1, 8, 0.3, "🔗 视频链接: ____________________", size=14, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# SLIDE: 亚洲在哪里
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🌏 亚洲在哪里？ Where is Asia?")
tb(s, 0.4, 0.9, 4.5, 0.5, "亚洲是七大洲中最大的洲！", size=20, bold=True, color=ORANGE)
tb(s, 0.4, 1.4, 4.5, 0.4, "Asia is the largest continent!", size=14, color=GRAY)
img_box(s, 5.2, 0.8, 4.5, 4.2, "📷 世界地图 World Map\n标出亚洲位置")
tb(s, 0.4, 2.0, 4.5, 0.4, "We are here! 我们在这里！ →", size=16, bold=True, color=DARK)
deco_corners(s, "🌏🐼🗻🐘")
page_num(s, n)

# ═══════════════════════════════════════════
# SLIDE 6: 认识亚洲 (1) - 地理
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🌏 认识亚洲  About Asia (1/2)")
items = [
    ("🏔️ 世界最大的洲 Largest Continent", "面积4,458万km² — 占地球30%"),
    ("🌐 48个国家 48 Countries", "亚洲有48个国家，人口46亿！"),
]
y = 0.9
for title, detail in items:
    tb(s, 0.4, y, 4.8, 0.5, title, size=18, bold=True, color=ORANGE)
    tb(s, 0.4, y+0.5, 4.8, 0.4, detail, size=15, color=DARK)
    y += 1.2
img_box(s, 5.5, 0.9, 4.2, 3.5, "📷 亚洲地图/风景图片")
page_num(s, n)

# ═══════════════════════════════════════════
# SLIDE 7: 认识亚洲 (2) - 特点
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🌏 认识亚洲  About Asia (2/2)")
items = [
    ("⛰️ 珠穆朗玛峰 Mt. Everest", "世界最高的山 — 8,849米！"),
    ("🗣️ 2,300+种语言", "亚洲是世界上语言最多的洲！"),
]
y = 0.9
for title, detail in items:
    tb(s, 0.4, y, 4.8, 0.5, title, size=18, bold=True, color=ORANGE)
    tb(s, 0.4, y+0.5, 4.8, 0.4, detail, size=15, color=DARK)
    y += 1.2
img_box(s, 5.5, 0.9, 4.2, 3.5, "📷 珠穆朗玛峰 / 亚洲多元文化")
page_num(s, n)

# SLIDE: 中国视频
s = new_slide(); n+=1
bg(s, RGBColor(0x7F,0x00,0x00))
tb(s, 1, 0.8, 8, 0.8, "🎬 看视频  Watch Video", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1, 1.8, 8, 0.5, "🇨🇳 认识中国 About China", size=22, color=RGBColor(0xFF,0xF3,0xE0), align=PP_ALIGN.CENTER)
img_box(s, 1.5, 2.5, 7, 2.5, "📷 插入视频截图或粘贴视频链接")
tb(s, 1, 5.1, 8, 0.3, "🔗 视频链接: ____________________", size=14, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# CHINA SLIDES
# ═══════════════════════════════════════════
# Slide: 中国 国旗+首都
s = country_info_slide("🇨🇳", "中国", "China", CHINA_RED, "国旗 + 首都",
    [("🏴 国旗 National Flag", "红色，五颗黄星 Red with 5 yellow stars"),
     ("🏛️ 首都 Capital", "北京 Beijing")],
    "📷 中国国旗 + 北京天安门"); n+=1; page_num(s, n)

# Slide: 中国 人口+语言
s = country_info_slide("🇨🇳", "中国", "China", CHINA_RED, "人口 + 语言",
    [("👥 人口 Population", "约14亿（世界第一！）~1.4 billion"),
     ("🗣️ 语言 Language", "中文（普通话）Chinese Mandarin")],
    "📷 中国人口/语言相关图片"); n+=1; page_num(s, n)

# Slide: 中国 主要景点
s = country_info_slide("🇨🇳", "中国", "China", CHINA_RED, "主要景点 Landmarks",
    [("🧱 长城 Great Wall", "21,196公里长！世界最长的城墙"),
     ("🏯 故宫 Forbidden City", "北京，有600年历史"),
     ("🐼 国宝：大熊猫 Giant Panda", "中国的国宝动物")],
    "📷 长城/故宫/熊猫 图片"); n+=1; page_num(s, n)

# Slide: 中国 礼节 Etiquette
s = country_info_slide("🇨🇳", "中国", "China", CHINA_RED, "礼节 Etiquette",
    [("👋 说你好 Say Hello", "「你好」nǐ hǎo"),
     ("🤝 打招呼 Greeting", "握手 Shake hands（不拥抱）"),
     ("🥢 吃饭礼节 Dining", "用筷子吃饭 / 筷子不能插在饭里（不吉利）")],
    "📷 握手 / 筷子 图片"); n+=1; page_num(s, n)

# Slide: 中国 美食 Food
s = country_info_slide("🇨🇳", "中国", "China", CHINA_RED, "美食 Food",
    [("🥟 饺子 Dumplings", "过年必吃！New Year must-eat!"),
     ("🍜 面条 Noodles / 炒饭 Fried Rice", "南方米饭，北方面食"),
     ("🦆 北京烤鸭 Peking Duck", "600年历史！600 years of history!")],
    "📷 中国美食图片"); n+=1; page_num(s, n)

# SLIDE: 日本视频
s = new_slide(); n+=1
bg(s, RGBColor(0x5D,0x00,0x17))
tb(s, 1, 0.8, 8, 0.8, "🎬 看视频  Watch Video", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1, 1.8, 8, 0.5, "🇯🇵 认识日本 About Japan", size=22, color=RGBColor(0xFF,0xF3,0xE0), align=PP_ALIGN.CENTER)
img_box(s, 1.5, 2.5, 7, 2.5, "📷 插入视频截图或粘贴视频链接")
tb(s, 1, 5.1, 8, 0.3, "🔗 视频链接: ____________________", size=14, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# JAPAN SLIDES
# ═══════════════════════════════════════════
s = country_info_slide("🇯🇵", "日本", "Japan", JAPAN_RED, "国旗 + 首都",
    [("🏴 国旗 National Flag", "白底红圆（太阳）White + Red circle"),
     ("🏛️ 首都 Capital", "东京 Tokyo")],
    "📷 日本国旗 + 东京塔"); n+=1; page_num(s, n)

s = country_info_slide("🇯🇵", "日本", "Japan", JAPAN_RED, "人口 + 语言",
    [("👥 人口 Population", "约1.25亿 ~125 million"),
     ("🗣️ 语言 Language", "日语 Japanese")],
    "📷 日本城市/文字图片"); n+=1; page_num(s, n)

s = country_info_slide("🇯🇵", "日本", "Japan", JAPAN_RED, "主要景点 Landmarks",
    [("🗻 富士山 Mt. Fuji", "日本最高峰 3,776米"),
     ("🌸 樱花 Cherry Blossom", "日本的象征，每年春天盛开"),
     ("🏝️ 6,852个岛屿！", "日本是一个岛国")],
    "📷 富士山/樱花 图片"); n+=1; page_num(s, n)

s = country_info_slide("🇯🇵", "日本", "Japan", JAPAN_RED, "礼节 Etiquette",
    [("👋 说你好 Say Hello", "「こんにちは」Konnichiwa (kon-ni-chi-wa)"),
     ("🙇 打招呼 Greeting", "鞠躬 Bow（弯腰越深越尊重，不握手！）"),
     ("🍜 吃饭礼节 Dining", "吃面可以发出声音（是礼貌！）/ 进屋脱鞋 / 不给小费")],
    "📷 鞠躬 / 吃拉面 图片"); n+=1; page_num(s, n)

# Slide: 日本 美食 Food
s = country_info_slide("🇯🇵", "日本", "Japan", JAPAN_RED, "美食 Food",
    [("🍣 寿司 Sushi", "全球最受欢迎的日本料理"),
     ("🍜 拉面 Ramen", "日本有5万多家拉面店！"),
     ("🍵 抹茶 Matcha / 🍱 便当 Bento", "茶道文化 / 午餐像艺术品！")],
    "📷 日本美食图片"); n+=1; page_num(s, n)

# SLIDE: 日本流行文化 Pop Culture
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🇯🇵 日本流行文化  Japanese Pop Culture", JAPAN_RED)
tb(s, 0.4, 0.85, 9, 0.35, "日本文化影响了全世界！These are all from Japan!", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# 2x3 grid of pop culture items with image placeholders
pop_items = [
    ("⚡ Pokemon 宝可梦", "全世界最赚钱的IP！\nPikachu你认识吗？"),
    ("🎵 Nintendo 任天堂", "Mario 马里奥\nSwitch 游戏机"),
    ("🐉 Dragon Ball 龙珠", "孙悟空！经典动漫\n影响了全世界"),
]
for i, (title, desc) in enumerate(pop_items):
    x = 0.3 + i * 3.15
    y = 1.2
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.95), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = [BG_WARM, RGBColor(0xE3,0xF2,0xFD), RGBColor(0xFC,0xE4,0xEC)][i]
    bar.line.fill.background()
    tb(s, x+0.1, y+0.05, 2.75, 0.4, title, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    lines = desc.split('\n')
    tf_d = tb(s, x+0.1, y+0.45, 2.75, 0.5, lines[0], size=12, color=DARK, align=PP_ALIGN.CENTER)
    for line in lines[1:]:
        add_p(tf_d, line, size=12, color=DARK, align=PP_ALIGN.CENTER)
# Large image/video area below
img_box(s, 0.3, 2.5, 9.4, 2.7, "📷 插入图片或视频 Insert images or video here")

deco_corners(s, "🌸🎮")
page_num(s, n)

# SLIDE: 印度视频
s = new_slide(); n+=1
bg(s, RGBColor(0x0A,0x44,0x04))
tb(s, 1, 0.8, 8, 0.8, "🎬 看视频  Watch Video", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1, 1.8, 8, 0.5, "🇮🇳 认识印度 About India", size=22, color=RGBColor(0xFF,0xF3,0xE0), align=PP_ALIGN.CENTER)
img_box(s, 1.5, 2.5, 7, 2.5, "📷 插入视频截图或粘贴视频链接")
tb(s, 1, 5.1, 8, 0.3, "🔗 视频链接: ____________________", size=14, color=RGBColor(0xBB,0xBB,0xBB), align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# INDIA SLIDES
# ═══════════════════════════════════════════
s = country_info_slide("🇮🇳", "印度", "India", INDIA_GREEN, "国旗 + 首都",
    [("🏴 国旗 National Flag", "橙白绿三色 + 蓝色法轮 Saffron/White/Green + Chakra"),
     ("🏛️ 首都 Capital", "新德里 New Delhi")],
    "📷 印度国旗 + 新德里"); n+=1; page_num(s, n)

s = country_info_slide("🇮🇳", "印度", "India", INDIA_GREEN, "人口 + 语言",
    [("👥 人口 Population", "约14亿（世界第二！）~1.4 billion"),
     ("🗣️ 语言 Language", "印地语+英语（22种官方语言！）")],
    "📷 印度人口/语言图片"); n+=1; page_num(s, n)

s = country_info_slide("🇮🇳", "印度", "India", INDIA_GREEN, "主要景点 Landmarks",
    [("🕌 泰姬陵 Taj Mahal", "世界七大奇迹之一！"),
     ("🏞️ 恒河 Ganges River", "印度最神圣的河流"),
     ("🐅 国家动物：孟加拉虎", "Bengal Tiger")],
    "📷 泰姬陵/恒河 图片"); n+=1; page_num(s, n)

s = country_info_slide("🇮🇳", "印度", "India", INDIA_GREEN, "礼节 Etiquette",
    [("👋 说你好 Say Hello", "「Namaste」नमस्ते（那马斯特）"),
     ("🙏 打招呼 Greeting", "双手合十 Palms together（不握手）"),
     ("🤚 吃饭礼节 Dining", "用右手吃饭（左手不干净）/ 很多人吃素 / 不吃牛肉（牛是神圣的）")],
    "📷 合十 / 用手吃饭 图片"); n+=1; page_num(s, n)

# Slide: 印度 美食 Food
s = country_info_slide("🇮🇳", "印度", "India", INDIA_GREEN, "美食 Food",
    [("🍛 咖喱 Curry", "印度最有名的食物，30+种咖喱！"),
     ("🫓 飞饼 Naan", "用来蘸咖喱吃 Dip in curry!"),
     ("☕ 奶茶 Chai", "印度国民饮料，每天10亿杯！")],
    "📷 印度美食图片"); n+=1; page_num(s, n)

# ═══════════════════════════════════════════
# SLIDE: 三国对比 Mini Role Play
# ═══════════════════════════════════════════
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🎭 打招呼练习  Greeting Practice")
countries = [
    ("🇨🇳 中国", "握手 +「你好！」", CHINA_RED),
    ("🇯🇵 日本", "鞠躬 +「こんにちは」(kon-ni-chi-wa)", JAPAN_RED),
    ("🇮🇳 印度", "合十 +「Namaste」", INDIA_GREEN),
]
for i, (name, greeting, clr) in enumerate(countries):
    x = 0.4 + i * 3.2
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.0), Inches(2.9), Inches(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bar.line.color.rgb = clr
    bar.line.width = Pt(3)
    tb(s, x+0.1, 1.1, 2.7, 0.5, name, size=22, bold=True, color=clr, align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 1.7, 2.7, 0.5, greeting, size=16, color=DARK, align=PP_ALIGN.CENTER)
    img_box(s, x+0.3, 2.4, 2.3, 1.8, "📷 动作示范")
tb(s, 0.4, 4.7, 9, 0.4, "站起来和旁边的同学练习！ Stand up and practice!", size=14, color=GRAY, align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# SESSION 2 DIVIDER
# ═══════════════════════════════════════════
section_divider("Session 2  下午", "复习总结 + 语言目标\n我会认：亚洲 中国 日本 印度 首都\n我会写：中国 日本 亚洲 首都", RGBColor(0x19,0x76,0xD2), "📖"); n+=1

# SLIDE: 快速复习
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "📖 快速复习  Quick Review", RGBColor(0x19,0x76,0xD2))
for i, (flag, name, q) in enumerate([
    ("🇨🇳", "中国", '怎么说"你好"？___'),
    ("🇯🇵", "日本", '怎么说"你好"？___'),
    ("🇮🇳", "印度", '怎么说"你好"？___'),
]):
    x = 0.4 + i * 3.2
    tb(s, x, 1.0, 2.8, 0.5, f"{flag} {name}", size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, x, 1.6, 2.8, 0.5, q, size=16, color=GRAY, align=PP_ALIGN.CENTER)
    img_box(s, x+0.2, 2.3, 2.4, 2.0, "📷")
tb(s, 0.4, 4.6, 9, 0.4, "谁能第一个说出来？ Who can answer first?", size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
page_num(s, n)

# SLIDE: 文化对比表 — 空白版（让学生回答）
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🌏 亚洲三国文化对比  你知道吗？", RGBColor(0x19,0x76,0xD2))
tb(s, 0.4, 0.85, 9, 0.35, "你能填出来吗？ Can you fill in the blanks?", size=14, color=GRAY, align=PP_ALIGN.CENTER)

table_shape_blank = s.shapes.add_table(7, 4, Inches(0.3), Inches(1.25), Inches(9.4), Inches(3.9))
tbl_b = table_shape_blank.table
tbl_b.columns[0].width = Inches(1.8)
tbl_b.columns[1].width = Inches(2.5)
tbl_b.columns[2].width = Inches(2.5)
tbl_b.columns[3].width = Inches(2.6)

data_blank = [
    ["", "🇨🇳 中国 China", "🇯🇵 日本 Japan", "🇮🇳 印度 India"],
    ["👋 说你好\nSay Hello", "？", "？", "？"],
    ["🤝 打招呼\nGreeting", "？", "？", "？"],
    ["🍽️ 吃饭工具\nEat with", "？", "？", "？"],
    ["⚠️ 吃饭礼节\nDining rule", "？", "？", "？"],
    ["🏛️ 首都\nCapital", "？", "？", "？"],
    ["🐾 代表动物\nAnimal", "？", "？", "？"],
]
for r, row_data in enumerate(data_blank):
    for c, cell_text in enumerate(row_data):
        cell = tbl_b.cell(r, c)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = cell_text
        run.font.size = Pt(11) if r > 0 else Pt(13)
        run.font.bold = True if (r == 0 or c == 0) else False
        run.font.name = 'Noto Sans SC'
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x19,0x76,0xD2)
        elif c == 0:
            run.font.color.rgb = RGBColor(0x44,0x44,0x44)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WARM
        else:
            run.font.color.rgb = RGBColor(0xCC,0xCC,0xCC)
            run.font.size = Pt(20)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF5)
page_num(s, n)

# SLIDE: 文化对比表 — 答案版
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🌏 亚洲三国文化对比  Comparison", RGBColor(0x19,0x76,0xD2))
tb(s, 0.4, 0.85, 9, 0.35, "三个国家各有特色，你最想去哪个？", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Build comparison table
from pptx.util import Emu
table_shape = s.shapes.add_table(7, 4, Inches(0.3), Inches(1.25), Inches(9.4), Inches(3.9))
tbl = table_shape.table

# Column widths
tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(2.5)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(2.6)

data = [
    ["", "🇨🇳 中国 China", "🇯🇵 日本 Japan", "🇮🇳 印度 India"],
    ["👋 说你好\nSay Hello", "你好\nnǐ hǎo", "こんにちは\nKonnichiwa (kon-ni-chi-wa)", "Namaste\n🙏"],
    ["🤝 打招呼\nGreeting", "握手\nShake hands", "鞠躬\nBow", "双手合十\nPalms together"],
    ["🍽️ 吃饭工具\nEat with", "筷子\nChopsticks 🥢", "筷子\nChopsticks 🥢", "右手\nRight hand 🤚"],
    ["⚠️ 吃饭礼节\nDining rule", "筷子不插饭里\nNo sticking!", "吃面可出声\nSlurp = polite!", "不用左手\nNo left hand!"],
    ["🏛️ 首都\nCapital", "北京\nBeijing", "东京\nTokyo", "新德里\nNew Delhi"],
    ["🐾 代表动物\nAnimal", "🐼 熊猫\nPanda", "🌸 樱花/🦢 鹤\nCherry/Crane", "🐅 老虎\nTiger"],
]

for r, row_data in enumerate(data):
    for c, cell_text in enumerate(row_data):
        cell = tbl.cell(r, c)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = cell_text
        run.font.size = Pt(11) if r > 0 else Pt(13)
        run.font.bold = True if (r == 0 or c == 0) else False
        run.font.name = 'Noto Sans SC'
        run.font.color.rgb = WHITE if r == 0 else (DARK if c > 0 else RGBColor(0x44,0x44,0x44))
        # Header row color
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x19,0x76,0xD2)
        elif c == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WARM
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF5)

page_num(s, n)

# ═══════════════════════════════════════════
# 我会认 — 5 individual word card slides
# ═══════════════════════════════════════════
BLUE = RGBColor(0x19,0x76,0xD2)
word_cards = [
    ("亚洲", "yà zhōu", "Asia", "亚洲是世界上最大的洲。", "📷 亚洲地图/风景"),
    ("中国", "zhōng guó", "China", "中国的首都是北京。", "📷 中国长城/天安门"),
    ("日本", "rì běn", "Japan", "日本有很多美丽的樱花。", "📷 日本富士山/樱花"),
    ("印度", "yìn dù", "India", "在印度，人们说 Namaste。", "📷 印度泰姬陵"),
    ("首都", "shǒu dū", "Capital", "北京是中国的首都。", "📷 各国首都建筑"),
]
for word, pinyin, eng, sentence, img_label in word_cards:
    s = new_slide(); n+=1
    bg(s, BG_CREAM)
    header_bar(s, "👀 我会认  I Can Read", BLUE)
    # Big word card (left)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0), Inches(4.5), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_WARM
    bar.line.fill.background()
    tb(s, 0.5, 1.1, 4.3, 1.4, word, size=72, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s, 0.5, 2.4, 4.3, 0.4, f"{pinyin}  {eng}", size=20, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, 0.5, 2.85, 4.3, 0.4, f"👉 跟我读！Read after me!", size=14, color=BLUE, align=PP_ALIGN.CENTER)
    # Image placeholder (right)
    img_box(s, 5.3, 1.0, 4.4, 2.5, img_label)
    # Example sentence (bottom)
    bar2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.8), Inches(9.2), Inches(1.2))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = WHITE
    bar2.line.color.rgb = BLUE
    bar2.line.width = Pt(2)
    tb(s, 0.6, 3.9, 1.5, 0.4, "例句", size=16, bold=True, color=BLUE)
    tb(s, 0.6, 4.3, 8.8, 0.5, sentence, size=22, bold=True, color=DARK)
    page_num(s, n)

# SLIDE: 练一练 — 词语游戏选择
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🎮 练一练  Word Games (选一个玩！)", BLUE)
games = [
    ("1️⃣", "拍苍蝇 Fly Swatter", "把字卡贴在白板上\n老师说词语，学生用拍子拍！", RGBColor(0xFF,0xF3,0xE0)),
    ("2️⃣", "举牌游戏 Show Me", "每人5张字卡\n老师说词语，学生举起正确的卡", RGBColor(0xE3,0xF2,0xFD)),
    ("3️⃣", "抢椅子 Musical Chairs", "椅子上放字卡\n音乐停，读出椅子上的词", RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣", "传话筒 Pass the Mic", "传球/话筒，停下的人\n读一张字卡并造句", RGBColor(0xFC,0xE4,0xEC)),
]
for i, (num, name, desc, bg_c) in enumerate(games):
    x = 0.3 + i * 2.4
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.9), Inches(2.2), Inches(4.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg_c
    bar.line.fill.background()
    tb(s, x+0.1, 1.0, 2.0, 0.4, num, size=24, align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 1.4, 2.0, 0.6, name, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    lines = desc.split('\n')
    tf = tb(s, x+0.15, 2.1, 1.9, 1.5, lines[0], size=13, color=DARK, align=PP_ALIGN.CENTER)
    for line in lines[1:]:
        add_p(tf, line, size=13, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 3.8, 2.0, 0.4, "低prep ✅", size=12, bold=True, color=RGBColor(0x38,0x8E,0x3C), align=PP_ALIGN.CENTER)
tb(s, 0.4, 5.15, 9, 0.3, "所有游戏只需要字卡，不需要额外准备！ All games just need word cards!", size=12, color=GRAY, align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# 我会写 — 3 individual writing slides (亚洲、中国、日本)
# ═══════════════════════════════════════════
write_cards = [
    ("亚洲", "yà zhōu", "Asia", "📷 亚洲地图"),
    ("中国", "zhōng guó", "China", "📷 中国国旗/地图"),
    ("日本", "rì běn", "Japan", "📷 日本国旗/地图"),
]
for word, pinyin, eng, img_label in write_cards:
    s = new_slide(); n+=1
    bg(s, BG_CREAM)
    header_bar(s, "✍️ 我会写  I Can Write", BLUE)
    # Big word (left)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0), Inches(4.5), Inches(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.color.rgb = BLUE
    bar.line.width = Pt(3)
    tb(s, 0.5, 1.05, 4.3, 1.2, word, size=72, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tb(s, 0.5, 2.2, 4.3, 0.4, f"{pinyin}  {eng}", size=20, color=GRAY, align=PP_ALIGN.CENTER)
    # Image placeholder (right)
    img_box(s, 5.3, 1.0, 4.4, 2.0, img_label)
    # Stroke order image placeholder (bottom-left)
    bar2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.3), Inches(5.0), Inches(1.8))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    bar2.line.fill.background()
    tb(s, 0.6, 3.4, 4.6, 0.4, "📝 笔顺 Stroke Order", size=16, bold=True, color=BLUE)
    img_box(s, 0.6, 3.9, 4.6, 1.0, "📷 插入笔顺图片")
    # Practice steps (right side of bottom)
    tf = tb(s, 5.8, 3.4, 3.8, 0.4, "练习步骤 Practice:", size=14, bold=True, color=BLUE)
    add_p(tf, "1. 空中写 Air Write", size=13, color=DARK)
    add_p(tf, "2. 手心写 Palm Write", size=13, color=DARK)
    add_p(tf, "3. 纸上写 Write 3 times", size=13, color=DARK)
    page_num(s, n)

# ═══════════════════════════════════════════
# SESSION 3 DIVIDER
# ═══════════════════════════════════════════
section_divider("Session 3  下午", "写Booklet + 做Project\n3:00 - 4:30", RGBColor(0x38,0x8E,0x3C), "🎨"); n+=1

# SLIDE: 写Booklet
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, '📓 完成"探索亚洲"练习册', RGBColor(0x38,0x8E,0x3C))
img_box(s, 0.4, 0.9, 9.2, 4.3, "📷 练习册截图")
page_num(s, n)

# SLIDE: Project Overview (keep as summary)
s = new_slide(); n+=1
bg(s, BG_CREAM)
header_bar(s, "🎨 Project Time!  4个手工项目", RGBColor(0x38,0x8E,0x3C))
projects = [
    ("PROJECT 1", "🧩 亚洲拼图", "group project", BG_WARM),
    ("PROJECT 2", "🪭 折扇", "独立完成", RGBColor(0xE3,0xF2,0xFD)),
    ("PROJECT 3", "🦢 Origami", "独立完成", RGBColor(0xE8,0xF5,0xE9)),
    ("PROJECT 4", "👗 印度服装", "group project", RGBColor(0xFC,0xE4,0xEC)),
]
for i, (proj, name, mode, bg_c) in enumerate(projects):
    x = 0.3 + i * 2.4
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.9), Inches(2.2), Inches(4.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg_c
    bar.line.fill.background()
    tb(s, x+0.1, 1.0, 2.0, 0.3, proj, size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 1.3, 2.0, 0.6, name, size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 1.9, 2.0, 0.3, f"({mode})", size=11, color=GRAY, align=PP_ALIGN.CENTER)
    img_box(s, x+0.2, 2.4, 1.8, 2.3, "📷 示范")
page_num(s, n)

# 4 individual project slides
proj_details = [
    ("PROJECT 1", "🧩 亚洲拼图  Asia Puzzle Map", "group project", BG_WARM, RGBColor(0x38,0x8E,0x3C)),
    ("PROJECT 2", "🪭 折扇  Chinese Folding Fan", "独立完成", RGBColor(0xE3,0xF2,0xFD), RGBColor(0x19,0x76,0xD2)),
    ("PROJECT 3", "🦢 Origami  日本折纸", "独立完成", RGBColor(0xE8,0xF5,0xE9), RGBColor(0x38,0x8E,0x3C)),
    ("PROJECT 4", "👗 印度传统服装  Indian Sari", "group project", RGBColor(0xFC,0xE4,0xEC), RGBColor(0xC2,0x18,0x5B)),
]
for proj, name, mode, bg_c, clr in proj_details:
    s = new_slide(); n+=1
    bg(s, bg_c)
    header_bar(s, f"{proj}: {name}", clr)
    tb(s, 0.4, 0.85, 9, 0.3, f"({mode})", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    # Left: image/screenshot placeholder
    img_box(s, 0.4, 1.3, 4.4, 3.5, "📷 示范图片 / 截图")
    # Right: video placeholder
    img_box(s, 5.2, 1.3, 4.5, 3.5, "🎬 教学视频 / 步骤视频")
    page_num(s, n)

# SLIDE: 签证章
s = new_slide(); n+=1
bg(s, BG_CREAM)
tb(s, 1, 0.5, 8, 0.8, "🪪 亚洲签证章  Asia Visa Stamp", size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
shape = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(1.5), Inches(3), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = ORANGE
shape.line.width = Pt(5)
tf = tb(s, 3.6, 1.8, 2.8, 2.5, "ASIA\n亚洲", size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_p(tf, "✓ VISITED", size=16, bold=True, color=RGBColor(0x38,0x8E,0x3C), align=PP_ALIGN.CENTER)
add_p(tf, "6/8/2025", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_p(tf, "中国 · 日本 · 印度", size=12, color=DARK, align=PP_ALIGN.CENTER)
tb(s, 1, 4.7, 8, 0.4, "恭喜你完成亚洲之旅！Congratulations! 🎉", size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
page_num(s, n)

# SLIDE: 明天航班
s = new_slide(); n+=1
bg(s, ORANGE)
tb(s, 1, 1.0, 8, 0.8, "✈️ 明天航班  Tomorrow's Flight", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tf = tb(s, 2, 2.2, 6, 2.5, "Flight 航班: GR-002", size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_p(tf, "Destination 目的地: 非洲 AFRICA 🌍", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_p(tf, "", size=10)
add_p(tf, "明天我们去非洲！", size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_p(tf, "那里的人怎么打招呼？", size=18, color=RGBColor(0xFF,0xF3,0xE0), align=PP_ALIGN.CENTER)
add_p(tf, "", size=10)
add_p(tf, "See you tomorrow, explorers! 明天见！", size=16, color=RGBColor(0xFF,0xF3,0xE0), align=PP_ALIGN.CENTER)
page_num(s, n)

# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
OUT = '/Users/Huan/projects/summercourse/Chinese/世界旅行world_trip_pbl/day1_asia_v2.pptx'
prs.save(OUT)
print(f"Created {n} slides → {OUT}")
