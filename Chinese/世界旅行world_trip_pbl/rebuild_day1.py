#!/usr/bin/env python3
"""
Rebuild day1_asia.pptx with new course structure.
Approach: Add new slides to the OLD file, reorder to desired sequence,
and keep only the slides we want (orphaned slides just stay in zip but don't show).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SRC = '/Users/Huan/Downloads/day1_asia-OLD.pptx'
OUT = '/Users/Huan/projects/summercourse/Chinese/世界旅行world_trip_pbl/day1_asia.pptx'

prs = Presentation(SRC)
slide_width = prs.slide_width
slide_height = prs.slide_height

# --- Helper functions ---
def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(textbox, text, size=Pt(18), bold=False, color=None, alignment=None, font_name='Noto Sans SC'):
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return tf

def add_paragraph(tf, text, size=Pt(18), bold=False, color=None, alignment=None, font_name='Noto Sans SC'):
    p = tf.add_paragraph()
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return p

def add_shape_bg(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_full_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape

def create_blank_slide(prs):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide

# --- Colors ---
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
DARK_ORANGE = RGBColor(0xE6, 0x7E, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
BLUE = RGBColor(0x19, 0x76, 0xD2)
GREEN = RGBColor(0x38, 0x8E, 0x3C)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
PINK = RGBColor(0xC2, 0x18, 0x5B)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
CARD_BG  = RGBColor(0xFF, 0xF3, 0xE0)
CARD_BG2 = RGBColor(0xE3, 0xF2, 0xFD)
CARD_BG3 = RGBColor(0xE8, 0xF5, 0xE9)
CARD_BG4 = RGBColor(0xFC, 0xE4, 0xEC)

# OLD file has 51 slides (indices 0-50)
# Record the index before adding new slides
old_count = len(prs.slides)
print(f"Original slide count: {old_count}")

# =============================================
# CREATE NEW SLIDES (added at end, indices 51+)
# =============================================

# --- SLIDE N0: 下午第一节 Divider ---
s = create_blank_slide(prs)
add_full_bg(s, ORANGE)
tb = add_textbox(s, Inches(1), Inches(0.8), Inches(8), Inches(1))
tf = set_text(tb, "☀\ufe0f 下午第一节", size=Pt(40), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(s, Inches(1), Inches(2.0), Inches(8), Inches(0.8))
set_text(tb2, "复习 → 认字活动 → 写字 → 完成Booklet", size=Pt(24), color=WHITE, alignment=PP_ALIGN.CENTER)
tb3 = add_textbox(s, Inches(2), Inches(3.2), Inches(6), Inches(1.5))
tf3 = set_text(tb3, "⏰ 50分钟", size=Pt(28), bold=True, color=RGBColor(0xFF,0xF3,0xE0), alignment=PP_ALIGN.CENTER)
add_paragraph(tf3, "复习Day 1 Lesson 1 → 认读生字 → 写字练习 → 完成booklet", size=Pt(18), color=WHITE, alignment=PP_ALIGN.CENTER)
N0 = old_count  # index 51

# --- SLIDE N1: 今日生字 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
set_text(tb, "📝 教学目标  Learning Objectives", size=Pt(30), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
# Row: 我会认
tb_obj1 = add_textbox(s, Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
set_text(tb_obj1, "👀 我会认 I can read:", size=Pt(20), bold=True, color=DARK_ORANGE)
words_row1 = [("亚洲", "yà zhōu", "Asia"), ("中国", "zhōng guó", "China"), ("日本", "rì běn", "Japan"), ("印度", "yìn dù", "India"), ("首都", "shǒu dū", "Capital")]
for i, (word, pinyin, eng) in enumerate(words_row1):
    x = Inches(0.3 + i * 1.9)
    add_shape_bg(s, x, Inches(1.4), Inches(1.7), Inches(1.6), CARD_BG)
    tb_c = add_textbox(s, x+Inches(0.05), Inches(1.45), Inches(1.6), Inches(0.8))
    set_text(tb_c, word, size=Pt(36), bold=True, color=DARK_ORANGE, alignment=PP_ALIGN.CENTER)
    tb_p = add_textbox(s, x+Inches(0.05), Inches(2.2), Inches(1.6), Inches(0.6))
    tf_p = set_text(tb_p, pinyin, size=Pt(14), color=BLACK, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf_p, eng, size=Pt(11), color=GRAY, alignment=PP_ALIGN.CENTER)
# Row: 我会写
tb_obj2 = add_textbox(s, Inches(0.5), Inches(3.2), Inches(9), Inches(0.5))
set_text(tb_obj2, "✍️ 我会写 I can write:", size=Pt(20), bold=True, color=BLUE)
words_row2 = [("中国", "zhōng guó"), ("日本", "rì běn"), ("亚洲", "yà zhōu"), ("首都", "shǒu dū")]
for i, (word, pinyin) in enumerate(words_row2):
    x = Inches(0.5 + i * 2.3)
    add_shape_bg(s, x, Inches(3.7), Inches(2.0), Inches(1.4), CARD_BG2)
    tb_c = add_textbox(s, x+Inches(0.1), Inches(3.75), Inches(1.8), Inches(0.8))
    set_text(tb_c, word, size=Pt(36), bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    tb_p = add_textbox(s, x+Inches(0.1), Inches(4.5), Inches(1.8), Inches(0.4))
    set_text(tb_p, pinyin, size=Pt(14), color=BLACK, alignment=PP_ALIGN.CENTER)
N1 = old_count + 1

# --- SLIDE N2: 闪卡认读 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "🎴 活动1：闪卡认读  Flashcard Reading", size=Pt(28), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
tb_l = add_textbox(s, Inches(0.5), Inches(1.2), Inches(4.5), Inches(3.8))
tf_l = set_text(tb_l, "玩法 How to Play:", size=Pt(20), bold=True, color=DARK)
add_paragraph(tf_l, "", size=Pt(6))
add_paragraph(tf_l, "1️⃣  老师举字卡，学生抢答", size=Pt(16), color=BLACK)
add_paragraph(tf_l, "     Teacher shows card, students race to answer", size=Pt(12), color=GRAY)
add_paragraph(tf_l, "", size=Pt(6))
add_paragraph(tf_l, "2️⃣  读出拼音 + 说出意思", size=Pt(16), color=BLACK)
add_paragraph(tf_l, "     Read pinyin + say the meaning", size=Pt(12), color=GRAY)
add_paragraph(tf_l, "", size=Pt(6))
add_paragraph(tf_l, "3️⃣  用这个字组词或造句", size=Pt(16), color=BLACK)
add_paragraph(tf_l, "     Make a word or sentence with it", size=Pt(12), color=GRAY)
add_paragraph(tf_l, "", size=Pt(6))
add_paragraph(tf_l, "4️⃣  答对得星星 ⭐", size=Pt(16), color=BLACK)
cards_data = [("亚洲", "yà zhōu", CARD_BG, DARK_ORANGE), ("中国", "zhōng guó", CARD_BG2, BLUE),
              ("日本", "rì běn", CARD_BG3, GREEN), ("印度", "yìn dù", CARD_BG4, PURPLE)]
for i, (txt, py, bg, clr) in enumerate(cards_data):
    x = Inches(5.5 + (i%2)*2.2)
    y = Inches(1.2 + (i//2)*1.8)
    add_shape_bg(s, x, y, Inches(1.8), Inches(1.6), bg)
    tb_c = add_textbox(s, x+Inches(0.1), y+Inches(0.1), Inches(1.6), Inches(1.0))
    set_text(tb_c, txt, size=Pt(36), bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    tb_p = add_textbox(s, x+Inches(0.1), y+Inches(1.0), Inches(1.6), Inches(0.4))
    set_text(tb_p, py, size=Pt(14), color=BLACK, alignment=PP_ALIGN.CENTER)
# 5th card: 首都
add_shape_bg(s, Inches(6.15), Inches(3.8), Inches(1.8), Inches(1.2), CARD_BG)
tb_5 = add_textbox(s, Inches(6.25), Inches(3.85), Inches(1.6), Inches(0.7))
set_text(tb_5, "首都", size=Pt(32), bold=True, color=RGBColor(0xE6,0x51,0x00), alignment=PP_ALIGN.CENTER)
tb_5p = add_textbox(s, Inches(6.25), Inches(4.5), Inches(1.6), Inches(0.4))
set_text(tb_5p, "shǒu dū", size=Pt(14), color=BLACK, alignment=PP_ALIGN.CENTER)
N2 = old_count + 2

# --- SLIDE N3: 字词配对 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "🔗 活动2：字词配对  Match Characters to Pictures", size=Pt(26), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
tb_i = add_textbox(s, Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
set_text(tb_i, "把汉字和图片连起来！Draw lines to match!", size=Pt(16), color=BLACK, alignment=PP_ALIGN.CENTER)
matches = [("饺子", "jiǎo zi", "🥟 dumplings"), ("寿司", "shòu sī", "🍣 sushi"),
           ("咖喱", "gā lí", "🍛 curry"), ("筷子", "kuài zi", "🥢 chopsticks"),
           ("面条", "miàn tiáo", "🍜 noodles")]
shuffled = [2, 4, 0, 3, 1]
for i, (char, pinyin, eng) in enumerate(matches):
    y = Inches(1.6 + i * 0.7)
    add_shape_bg(s, Inches(0.8), y, Inches(2.6), Inches(0.55), CARD_BG)
    tb_c = add_textbox(s, Inches(0.9), y+Inches(0.03), Inches(2.4), Inches(0.5))
    set_text(tb_c, f"{char}  ({pinyin})", size=Pt(16), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    j = shuffled[i]
    add_shape_bg(s, Inches(6.5), y, Inches(2.5), Inches(0.55), CARD_BG2)
    tb_e = add_textbox(s, Inches(6.6), y+Inches(0.03), Inches(2.3), Inches(0.5))
    set_text(tb_e, matches[j][2], size=Pt(16), color=DARK, alignment=PP_ALIGN.CENTER)
tb_a = add_textbox(s, Inches(3.5), Inches(2.5), Inches(3), Inches(0.5))
set_text(tb_a, "← 连线 Draw lines →", size=Pt(14), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
N3 = old_count + 3

# --- SLIDE N4: 写字练习 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "✍\ufe0f 我会写  I Can Write", size=Pt(30), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
practice_words = [("中国", "zhōng guó"), ("日本", "rì běn"), ("亚洲", "yà zhōu"), ("首都", "shǒu dū")]
for i, (word, pinyin) in enumerate(practice_words):
    x = Inches(0.5 + i * 2.4)
    add_shape_bg(s, x, Inches(1.2), Inches(2.1), Inches(1.5), WHITE)
    tb_c = add_textbox(s, x+Inches(0.1), Inches(1.2), Inches(1.9), Inches(1.0))
    set_text(tb_c, word, size=Pt(48), bold=True, color=DARK_ORANGE, alignment=PP_ALIGN.CENTER)
    tb_p = add_textbox(s, x+Inches(0.1), Inches(2.15), Inches(1.9), Inches(0.4))
    set_text(tb_p, pinyin, size=Pt(14), color=BLACK, alignment=PP_ALIGN.CENTER)
tb_pr = add_textbox(s, Inches(0.5), Inches(3.6), Inches(9), Inches(1.5))
tf_pr = set_text(tb_pr, "练习步骤 Practice Steps:", size=Pt(18), bold=True, color=DARK)
add_paragraph(tf_pr, "1. 空中写 Air Write — 跟老师一起用手指在空中写", size=Pt(15), color=BLACK)
add_paragraph(tf_pr, "2. 手心写 Palm Write — 在手心上写一写", size=Pt(15), color=BLACK)
add_paragraph(tf_pr, "3. 纸上写 Paper Write — 在练习纸上写3遍", size=Pt(15), color=BLACK)
add_paragraph(tf_pr, "4. 同伴检查 Peer Check — 互相检查，给星星 ⭐", size=Pt(15), color=BLACK)
N4 = old_count + 4

# --- SLIDE N5: 美食生字 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "🍜 美食生字  Food Characters", size=Pt(30), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
foods = [("饺子", "jiǎo zi", "🥟", "中国过年必吃！"), ("寿司", "shòu sī", "🍣", "日本最有名的食物"),
         ("咖喱", "gā lí", "🍛", "印度国民美食"), ("筷子", "kuài zi", "🥢", "中国人吃饭的工具"),
         ("面条", "miàn tiáo", "🍜", "北方人的最爱"), ("拉面", "lā miàn", "🍜", "日本有5万家拉面店")]
for i, (char, pinyin, emoji, desc) in enumerate(foods):
    row, col = i // 3, i % 3
    x, y = Inches(0.4 + col * 3.1), Inches(1.2 + row * 2.1)
    bg = [CARD_BG, CARD_BG2, CARD_BG3][col]
    add_shape_bg(s, x, y, Inches(2.8), Inches(1.8), bg)
    tb_e = add_textbox(s, x+Inches(0.1), y+Inches(0.05), Inches(2.6), Inches(0.5))
    set_text(tb_e, emoji, size=Pt(28), alignment=PP_ALIGN.CENTER)
    tb_c = add_textbox(s, x+Inches(0.1), y+Inches(0.5), Inches(2.6), Inches(0.6))
    set_text(tb_c, char, size=Pt(28), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    tb_py = add_textbox(s, x+Inches(0.1), y+Inches(1.0), Inches(2.6), Inches(0.35))
    set_text(tb_py, pinyin, size=Pt(14), color=BLACK, alignment=PP_ALIGN.CENTER)
    tb_d = add_textbox(s, x+Inches(0.1), y+Inches(1.35), Inches(2.6), Inches(0.35))
    set_text(tb_d, desc, size=Pt(11), color=GRAY, alignment=PP_ALIGN.CENTER)
N5 = old_count + 5

# --- SLIDE N6: 完成Booklet ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "📓 完成Booklet  Fill in Your Booklet", size=Pt(30), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
steps = [("1", "写国家名", "在地图上标注：中国、日本、印度\nLabel: China, Japan, India"),
         ("2", "画国旗", "画出三个国家的国旗\nDraw the three country flags"),
         ("3", "写生字", "抄写今天学的生字各2遍\nCopy today's characters x2"),
         ("4", "写句子", "用「我喜欢吃___」写一个句子\nWrite a sentence using the pattern")]
colors_list = [CARD_BG, CARD_BG2, CARD_BG3, CARD_BG4]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + (i%2)*4.8)
    y = Inches(1.2 + (i//2)*2.0)
    add_shape_bg(s, x, y, Inches(4.3), Inches(1.7), colors_list[i])
    tb_n = add_textbox(s, x+Inches(0.2), y+Inches(0.15), Inches(0.5), Inches(0.5))
    set_text(tb_n, num, size=Pt(24), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    tb_t = add_textbox(s, x+Inches(0.7), y+Inches(0.15), Inches(3.3), Inches(0.45))
    set_text(tb_t, title, size=Pt(17), bold=True, color=DARK)
    lines = desc.split('\n')
    tb_d = add_textbox(s, x+Inches(0.7), y+Inches(0.6), Inches(3.3), Inches(0.9))
    tf_d = set_text(tb_d, lines[0], size=Pt(13), color=BLACK)
    if len(lines) > 1:
        add_paragraph(tf_d, lines[1], size=Pt(11), color=GRAY)
N6 = old_count + 6

# --- SLIDE N7: 下午第二节 Project Divider ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0x1B,0x5E,0x20))
tb = add_textbox(s, Inches(1), Inches(0.6), Inches(8), Inches(1))
set_text(tb, "🎨 下午第二节 — Project Time!", size=Pt(38), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(s, Inches(1), Inches(1.8), Inches(8), Inches(0.5))
set_text(tb2, "⏰ 50分钟  4个手工项目", size=Pt(22), color=RGBColor(0xC8,0xE6,0xC9), alignment=PP_ALIGN.CENTER)
projects = [("PROJECT 1", "🧩 亚洲拼图", "Asia Puzzle Map", CARD_BG),
            ("PROJECT 2", "🪭 折扇", "Chinese Folding Fan", CARD_BG2),
            ("PROJECT 3", "🦢 Origami", "Japanese Paper Folding", CARD_BG3),
            ("PROJECT 4", "👗 印度传统服装", "Indian Sari", CARD_BG4)]
for i, (proj, cn, en, bg) in enumerate(projects):
    x = Inches(0.5 + i*2.35)
    add_shape_bg(s, x, Inches(2.6), Inches(2.1), Inches(2.3), bg)
    tb_p = add_textbox(s, x+Inches(0.1), Inches(2.7), Inches(1.9), Inches(0.4))
    set_text(tb_p, proj, size=Pt(12), bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)
    tb_cn = add_textbox(s, x+Inches(0.1), Inches(3.1), Inches(1.9), Inches(0.8))
    set_text(tb_cn, cn, size=Pt(22), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    tb_en = add_textbox(s, x+Inches(0.1), Inches(4.0), Inches(1.9), Inches(0.6))
    set_text(tb_en, en, size=Pt(11), color=RGBColor(0x66,0x66,0x66), alignment=PP_ALIGN.CENTER)
N7 = old_count + 7

# --- SLIDE N8: PROJECT 1 亚洲拼图 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFF,0xFA,0xF0))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "🧩 PROJECT 1：亚洲拼图  Asia Puzzle Map", size=Pt(28), bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
tb_t = add_textbox(s, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
set_text(tb_t, "⏰ 约10分钟", size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)
tb_l = add_textbox(s, Inches(0.5), Inches(1.4), Inches(5.5), Inches(3.5))
tf_l = set_text(tb_l, "材料 Materials:", size=Pt(16), bold=True, color=DARK)
add_paragraph(tf_l, "• 亚洲地图拼图模板（预先打印裁好）", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "• 彩色笔 / 蜡笔", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "• 胶棒", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "", size=Pt(6))
add_paragraph(tf_l, "步骤 Steps:", size=Pt(16), bold=True, color=DARK)
add_paragraph(tf_l, "1. 在拼图块上找到中国、日本、印度", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "2. 用不同颜色给三个国家涂色", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "3. 写上国家名字（汉字+拼音）", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "4. 把拼图拼好，贴在底板上", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "5. 画上每个国家的标志（长城/富士山/泰姬陵）", size=Pt(14), color=BLACK)
add_shape_bg(s, Inches(6.3), Inches(1.4), Inches(3.2), Inches(3.5), RGBColor(0xE0,0xE0,0xE0))
tb_img = add_textbox(s, Inches(6.4), Inches(2.5), Inches(3.0), Inches(1.0))
set_text(tb_img, "📷 示范图片\nSample Photo", size=Pt(16), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
N8 = old_count + 8

# --- SLIDE N9: PROJECT 2 折扇 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xE3,0xF2,0xFD))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "🪭 PROJECT 2：折扇  Chinese Folding Fan", size=Pt(28), bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
tb_t = add_textbox(s, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
set_text(tb_t, "⏰ 约10分钟", size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)
fan_steps = [("1", "准备", "彩色纸+冰棍棒\n+胶水"),
             ("2", "折叠", "来回折叠\n像手风琴"),
             ("3", "装饰", "画龙、花\n或写汉字"),
             ("4", "粘合", "底部粘合\n贴冰棍棒"),
             ("5", "写字", "写一个\n今天学的字")]
for i, (num, title, desc) in enumerate(fan_steps):
    x = Inches(0.3 + i*1.9)
    bg = [CARD_BG, CARD_BG2, CARD_BG3, CARD_BG4, CARD_BG][i]
    add_shape_bg(s, x, Inches(1.5), Inches(1.7), Inches(3.2), bg)
    tb_n = add_textbox(s, x+Inches(0.1), Inches(1.6), Inches(1.5), Inches(0.45))
    set_text(tb_n, f"Step {num}", size=Pt(14), bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    tb_ti = add_textbox(s, x+Inches(0.1), Inches(2.0), Inches(1.5), Inches(0.4))
    set_text(tb_ti, title, size=Pt(18), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_shape_bg(s, x+Inches(0.15), Inches(2.5), Inches(1.4), Inches(1.1), RGBColor(0xE0,0xE0,0xE0))
    tb_ph = add_textbox(s, x+Inches(0.15), Inches(2.8), Inches(1.4), Inches(0.6))
    set_text(tb_ph, "📷", size=Pt(20), color=RGBColor(0xBB,0xBB,0xBB), alignment=PP_ALIGN.CENTER)
    lines = desc.split('\n')
    tb_d = add_textbox(s, x+Inches(0.1), Inches(3.7), Inches(1.5), Inches(0.7))
    tf_d = set_text(tb_d, lines[0], size=Pt(12), color=BLACK, alignment=PP_ALIGN.CENTER)
    if len(lines) > 1:
        add_paragraph(tf_d, lines[1], size=Pt(12), color=BLACK, alignment=PP_ALIGN.CENTER)
N9 = old_count + 9

# --- SLIDE N10: PROJECT 3 Origami ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xE8,0xF5,0xE9))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "🦢 PROJECT 3：Origami  日本折纸", size=Pt(28), bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
tb_t = add_textbox(s, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
set_text(tb_t, "⏰ 约15分钟", size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)
tb_l = add_textbox(s, Inches(0.5), Inches(1.4), Inches(4.5), Inches(1.5))
tf_l = set_text(tb_l, "折纸 Origami = 折 (ori, fold) + 纸 (gami, paper)", size=Pt(15), bold=True, color=GREEN)
add_paragraph(tf_l, "", size=Pt(4))
add_paragraph(tf_l, "日本折纸有1000多年的历史！", size=Pt(14), color=BLACK)
add_paragraph(tf_l, "Origami has over 1,000 years of history!", size=Pt(12), color=GRAY)
add_paragraph(tf_l, "", size=Pt(4))
add_paragraph(tf_l, "材料: 正方形折纸 (origami paper)", size=Pt(14), bold=True, color=DARK)
add_shape_bg(s, Inches(5.5), Inches(1.4), Inches(4.0), Inches(1.5), RGBColor(0xE0,0xE0,0xE0))
tb_img = add_textbox(s, Inches(5.6), Inches(1.7), Inches(3.8), Inches(0.8))
set_text(tb_img, "📷 折纸示范图片\nOrigami Demo Photos", size=Pt(14), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
tb_o = add_textbox(s, Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
set_text(tb_o, "选一个折 Choose one to fold:", size=Pt(16), bold=True, color=DARK)
options = [("🦢 千纸鹤", "Paper Crane", "和平的象征", "⭐⭐⭐"),
           ("🐸 青蛙", "Jumping Frog", "会跳的青蛙！", "⭐⭐"),
           ("🌷 郁金香", "Tulip", "简单又漂亮", "⭐"),
           ("❤\ufe0f 爱心", "Heart", "送给爸爸妈妈", "⭐")]
for i, (name, en, desc, diff) in enumerate(options):
    x = Inches(0.3 + i*2.4)
    bg = [CARD_BG, CARD_BG2, CARD_BG3, CARD_BG4][i]
    add_shape_bg(s, x, Inches(3.6), Inches(2.1), Inches(1.6), bg)
    tb_n = add_textbox(s, x+Inches(0.1), Inches(3.65), Inches(1.9), Inches(0.45))
    set_text(tb_n, name, size=Pt(16), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    tb_en = add_textbox(s, x+Inches(0.1), Inches(4.05), Inches(1.9), Inches(0.3))
    set_text(tb_en, en, size=Pt(12), color=RGBColor(0x66,0x66,0x66), alignment=PP_ALIGN.CENTER)
    tb_df = add_textbox(s, x+Inches(0.1), Inches(4.35), Inches(1.9), Inches(0.3))
    set_text(tb_df, f"难度: {diff}", size=Pt(10), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    tb_d = add_textbox(s, x+Inches(0.1), Inches(4.7), Inches(1.9), Inches(0.3))
    set_text(tb_d, desc, size=Pt(11), color=BLACK, alignment=PP_ALIGN.CENTER)
N10 = old_count + 10

# --- SLIDE N11: PROJECT 4 印度传统服装 ---
s = create_blank_slide(prs)
add_full_bg(s, RGBColor(0xFC,0xE4,0xEC))
tb = add_textbox(s, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
set_text(tb, "👗 PROJECT 4：印度传统服装  Indian Sari", size=Pt(28), bold=True, color=PINK, alignment=PP_ALIGN.CENTER)
tb_t = add_textbox(s, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
set_text(tb_t, "⏰ 约15分钟", size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)
tb_intro = add_textbox(s, Inches(0.5), Inches(1.3), Inches(4.5), Inches(1.2))
tf_i = set_text(tb_intro, "🇮🇳 纱丽 Sari — 印度最美的传统服装", size=Pt(16), bold=True, color=PINK)
add_paragraph(tf_i, "• 一块长长的布（约5-9米！）", size=Pt(14), color=BLACK)
add_paragraph(tf_i, "• 颜色鲜艳，有很多装饰", size=Pt(14), color=BLACK)
add_paragraph(tf_i, "• 红色=喜庆  黄色=吉祥  绿色=丰收", size=Pt(12), color=RGBColor(0x66,0x66,0x66))
tb_m = add_textbox(s, Inches(0.5), Inches(2.7), Inches(4.5), Inches(1.2))
tf_m = set_text(tb_m, "材料 Materials:", size=Pt(16), bold=True, color=DARK)
add_paragraph(tf_m, "• 彩色褶皱纸 (crepe paper) 多种颜色", size=Pt(14), color=BLACK)
add_paragraph(tf_m, "• 亮片、小珠子（装饰用）", size=Pt(14), color=BLACK)
add_paragraph(tf_m, "• 胶棒 / 双面胶", size=Pt(14), color=BLACK)
add_paragraph(tf_m, "• 纸板人偶模板（预先准备）", size=Pt(14), color=BLACK)
tb_s = add_textbox(s, Inches(0.5), Inches(4.1), Inches(9), Inches(1.2))
tf_s = set_text(tb_s, "步骤: 1.选颜色 → 2.裁长条 → 3.缠绕出纱丽 → 4.用亮片装饰 → 5.写上\"印度\"+\"Namaste\"", size=Pt(13), bold=True, color=DARK)
add_shape_bg(s, Inches(5.5), Inches(1.3), Inches(4.0), Inches(2.5), RGBColor(0xE0,0xE0,0xE0))
tb_img = add_textbox(s, Inches(5.6), Inches(2.0), Inches(3.8), Inches(0.8))
set_text(tb_img, "📷 示范图片\nSample Photos", size=Pt(14), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
N11 = old_count + 11

# ============================================================
# COUNTRY SPLIT SLIDES: 每个国家拆成2页
# ============================================================
IMG_PLACEHOLDER = RGBColor(0xE0, 0xE0, 0xE0)

def make_country_slide_a(prs, flag_emoji, country_cn, country_en, color,
                          flag_desc, capital_cn, capital_en, treasure_cn, treasure_en):
    """Page A: 国旗 + 首都 + 国宝 (with image placeholders)"""
    s = create_blank_slide(prs)
    add_full_bg(s, RGBColor(0xFF, 0xFA, 0xF0))
    tb = add_textbox(s, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    set_text(tb, f"{flag_emoji} {country_cn} {country_en}", size=Pt(30), bold=True, color=color, alignment=PP_ALIGN.CENTER)

    # 国旗 Flag section
    add_shape_bg(s, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.45), color)
    tb_h = add_textbox(s, Inches(0.4), Inches(1.02), Inches(4.3), Inches(0.4))
    set_text(tb_h, "🏴 国旗 National Flag", size=Pt(16), bold=True, color=WHITE)
    # Flag image placeholder
    add_shape_bg(s, Inches(0.3), Inches(1.55), Inches(4.5), Inches(2.0), IMG_PLACEHOLDER)
    tb_img = add_textbox(s, Inches(0.4), Inches(2.1), Inches(4.3), Inches(0.8))
    set_text(tb_img, f"📷 {flag_desc}", size=Pt(14), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 首都 Capital section
    add_shape_bg(s, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.45), color)
    tb_h2 = add_textbox(s, Inches(5.3), Inches(1.02), Inches(4.3), Inches(0.4))
    set_text(tb_h2, "🏛️ 首都 Capital", size=Pt(16), bold=True, color=WHITE)
    # Capital image placeholder
    add_shape_bg(s, Inches(5.2), Inches(1.55), Inches(4.5), Inches(2.0), IMG_PLACEHOLDER)
    tb_img2 = add_textbox(s, Inches(5.3), Inches(1.7), Inches(4.3), Inches(0.5))
    set_text(tb_img2, f"📷 {capital_cn} {capital_en}", size=Pt(14), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Capital name
    tb_cap = add_textbox(s, Inches(5.3), Inches(3.0), Inches(4.3), Inches(0.45))
    set_text(tb_cap, f"{capital_cn}  {capital_en}", size=Pt(20), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

    # 国宝 National Treasure section (bottom full width)
    add_shape_bg(s, Inches(0.3), Inches(3.7), Inches(9.4), Inches(0.45), color)
    tb_h3 = add_textbox(s, Inches(0.4), Inches(3.72), Inches(9.2), Inches(0.4))
    set_text(tb_h3, f"🐾 国宝 National Treasure: {treasure_cn} {treasure_en}", size=Pt(16), bold=True, color=WHITE)
    # Treasure image placeholder
    add_shape_bg(s, Inches(2.5), Inches(4.25), Inches(5.0), Inches(1.2), IMG_PLACEHOLDER)
    tb_img3 = add_textbox(s, Inches(2.6), Inches(4.5), Inches(4.8), Inches(0.5))
    set_text(tb_img3, f"📷 {treasure_cn}", size=Pt(14), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return s

def make_country_slide_b(prs, flag_emoji, country_cn, country_en, color,
                          pop_cn, pop_en, lang_cn, lang_en, area_cn, fun_fact):
    """Page B: 人口 + 语言 + 面积 (with image placeholders)"""
    s = create_blank_slide(prs)
    add_full_bg(s, RGBColor(0xFF, 0xFA, 0xF0))
    tb = add_textbox(s, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    set_text(tb, f"{flag_emoji} {country_cn} {country_en}", size=Pt(30), bold=True, color=color, alignment=PP_ALIGN.CENTER)

    # 人口 Population (left)
    add_shape_bg(s, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.45), color)
    tb_h = add_textbox(s, Inches(0.4), Inches(1.02), Inches(4.3), Inches(0.4))
    set_text(tb_h, "👥 人口 Population", size=Pt(16), bold=True, color=WHITE)
    add_shape_bg(s, Inches(0.3), Inches(1.55), Inches(4.5), Inches(1.6), IMG_PLACEHOLDER)
    tb_pop = add_textbox(s, Inches(0.4), Inches(1.7), Inches(4.3), Inches(0.8))
    tf_pop = set_text(tb_pop, pop_cn, size=Pt(22), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf_pop, pop_en, size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)

    # 语言 Language (right)
    add_shape_bg(s, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.45), color)
    tb_h2 = add_textbox(s, Inches(5.3), Inches(1.02), Inches(4.3), Inches(0.4))
    set_text(tb_h2, "🗣️ 语言 Language", size=Pt(16), bold=True, color=WHITE)
    add_shape_bg(s, Inches(5.2), Inches(1.55), Inches(4.5), Inches(1.6), IMG_PLACEHOLDER)
    tb_lang = add_textbox(s, Inches(5.3), Inches(1.7), Inches(4.3), Inches(0.8))
    tf_lang = set_text(tb_lang, lang_cn, size=Pt(22), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf_lang, lang_en, size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)

    # 面积 / 有趣事实 (bottom)
    add_shape_bg(s, Inches(0.3), Inches(3.4), Inches(9.4), Inches(0.45), color)
    tb_h3 = add_textbox(s, Inches(0.4), Inches(3.42), Inches(9.2), Inches(0.4))
    set_text(tb_h3, f"📏 {area_cn}", size=Pt(16), bold=True, color=WHITE)

    # Fun facts row
    add_shape_bg(s, Inches(0.3), Inches(4.0), Inches(9.4), Inches(1.3), RGBColor(0xFF, 0xF3, 0xE0))
    tb_fun = add_textbox(s, Inches(0.5), Inches(4.1), Inches(9.0), Inches(1.1))
    tf_fun = set_text(tb_fun, "💡 你知道吗？Did you know?", size=Pt(18), bold=True, color=color)
    add_paragraph(tf_fun, fun_fact, size=Pt(16), color=DARK)
    return s

# --- 中国 Slide A: 国旗+首都+国宝 ---
make_country_slide_a(prs, "🇨🇳", "中国", "China", RGBColor(0xDE, 0x29, 0x10),
    "红色 + 五颗黄星 Red with 5 yellow stars",
    "北京", "Beijing", "大熊猫", "Giant Panda 🐼")
CN_A = old_count + 12

# --- 中国 Slide B: 人口+语言 ---
make_country_slide_b(prs, "🇨🇳", "中国", "China", RGBColor(0xDE, 0x29, 0x10),
    "约14亿（世界第一！）", "~1.4 billion (No.1!)",
    "中文（普通话）", "Chinese (Mandarin)",
    "面积：世界第三大  Covers 9.6M km²",
    "中国有五千年的历史！四大发明：造纸、印刷、火药、指南针\nChina has 5,000 years of history! Four Great Inventions!")
CN_B = old_count + 13

# --- 日本 Slide A: 国旗+首都+国宝 ---
make_country_slide_a(prs, "🇯🇵", "日本", "Japan", RGBColor(0xBC, 0x00, 0x2D),
    "白底红圆（太阳）White + Red circle",
    "东京", "Tokyo", "富士山 + 樱花", "Mt. Fuji + Cherry Blossom 🌸")
JP_A = old_count + 14

# --- 日本 Slide B: 人口+语言 ---
make_country_slide_b(prs, "🇯🇵", "日本", "Japan", RGBColor(0xBC, 0x00, 0x2D),
    "约1.25亿", "~125 million",
    "日语", "Japanese",
    "由6,852个岛屿组成！Made of 6,852 islands!",
    "日本是世界上最安全、最干净的国家之一！新干线时速320km!\nJapan is one of the safest & cleanest countries! Bullet train: 320km/h!")
JP_B = old_count + 15

# --- 印度 Slide A: 国旗+首都+国宝 ---
make_country_slide_a(prs, "🇮🇳", "印度", "India", RGBColor(0x13, 0x88, 0x08),
    "橙白绿三色 + 蓝色法轮 Saffron/White/Green + Blue Chakra",
    "新德里", "New Delhi", "孟加拉虎", "Bengal Tiger 🐅")
IN_A = old_count + 16

# --- 印度 Slide B: 人口+语言 ---
make_country_slide_b(prs, "🇮🇳", "印度", "India", RGBColor(0x13, 0x88, 0x08),
    "约14亿（世界第二！）", "~1.4 billion (No.2!)",
    "印地语 + 英语\n（22种官方语言！）", "Hindi + English\n(22 official languages!)",
    "面积：世界第七大  7th largest country",
    "印度发明了数字0！宝莱坞是世界最大电影业！\nIndia invented the number 0! Bollywood is the world's largest film industry!")
IN_B = old_count + 17

print(f"Total slides after adding new: {len(prs.slides)}")

# =============================================
# REORDER: Build the final slide sequence
# =============================================
# OLD file indices (0-50):
# 0: Boarding pass
# 1: 时间安排
# 2: Day 1 Lesson1
# 3: Weekly itinerary
# 4: Today's goals
# 5: 七大洲
# 6: 大洲排序
# 7: 亚洲在哪里 (map+text)
# 8: 亚洲地图 (image only)
# 9: 认识亚洲 facts
# 10: 中国在哪里
# 11: 看视频
# 12: 中国概览
# 13: 中国文化礼节
# 14: 中国美食
# 15: 练一练 Kahoot
# 16: (blank - was empty)
# 17-22: (blank/image slides)
# 23: 中国 Check 1/2
# 24: 中国 Check 2/2
# 25: 日本概览
# 26: 日本文化礼节
# 27: 日本美食
# 28: 日本 Check 1/2
# 29: 日本 Check 2/2
# 30: 印度概览
# 31: 印度文化礼节
# 32: 印度美食
# 33: 印度 Check 1/2
# 34: 印度 Check 2/2
# 35: Mini Role Play
# 36: Project 提醒
# 37: ☀️ 下午开始 (old divider)
# 38: Quick Review
# 39: Comparison Table
# 40: Similarities
# 41: Travel Tips
# 42: Vocab Cards
# 43: Sentence Patterns
# 44: Role Play
# 45: Project Passport
# 46: Project Differentiation
# 47: Project Example
# 48: Sharing Time
# 49: Visa Stamp
# 50: Tomorrow Flight

final_order = [
    # === 上午 Lesson 1 ===
    0,   # Boarding pass
    1,   # 时间安排
    2,   # Day 1 Lesson1
    3,   # Weekly itinerary
    4,   # Today's goals
    5,   # 七大洲
    6,   # 大洲排序
    7,   # 亚洲在哪里
    8,   # 亚洲地图
    9,   # 认识亚洲 facts
    10,  # 中国在哪里
    11,  # 看视频
    CN_A, # 中国: 国旗+首都+国宝
    CN_B, # 中国: 人口+语言+面积
    13,  # 中国文化礼节
    14,  # 中国美食
    15,  # 练一练 Kahoot
    23,  # 中国 Check 1
    24,  # 中国 Check 2
    JP_A, # 日本: 国旗+首都+国宝
    JP_B, # 日本: 人口+语言
    26,  # 日本文化礼节
    27,  # 日本美食
    28,  # 日本 Check 1
    29,  # 日本 Check 2
    IN_A, # 印度: 国旗+首都+国宝
    IN_B, # 印度: 人口+语言
    31,  # 印度文化礼节
    32,  # 印度美食
    33,  # 印度 Check 1
    34,  # 印度 Check 2
    35,  # Mini Role Play
    36,  # Project 提醒

    # === 下午第一节: 复习+认字+booklet ===
    N0,  # 下午第一节 divider
    38,  # Quick Review
    39,  # Comparison Table
    40,  # Similarities
    41,  # Travel Tips
    42,  # Vocab Cards
    N1,  # 今日生字
    N2,  # 闪卡认读
    N3,  # 字词配对
    N5,  # 美食生字
    N4,  # 写字练习
    43,  # Sentence Patterns
    44,  # Role Play
    N6,  # 完成Booklet

    # === 下午第二节: Projects ===
    N7,  # Project Time divider
    N8,  # Project 1 亚洲拼图
    N9,  # Project 2 折扇
    N10, # Project 3 Origami
    N11, # Project 4 印度传统服装
    45,  # Project Passport
    46,  # Project Differentiation
    47,  # Project Example
    48,  # Sharing Time
    49,  # Visa Stamp
    50,  # Tomorrow Flight
]

# Reorder slides
sldIdLst = prs.slides._sldIdLst
sldIds = list(sldIdLst)

new_order = [sldIds[i] for i in final_order]

for el in list(sldIdLst):
    sldIdLst.remove(el)
for el in new_order:
    sldIdLst.append(el)

print(f"Final presentation: {len(list(sldIdLst))} slides")

prs.save(OUT)
print(f"Saved to {OUT}")
