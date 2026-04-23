#!/usr/bin/env python3
"""
Day 2 Africa PPT — Same structure as Day 1 Asia v2.
Countries: Egypt 埃及, Kenya 肯尼亚, South Africa 南非
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# Colors
GOLD = RGBColor(0xD4,0xA0,0x17)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BLACK = RGBColor(0x33,0x33,0x33)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
EGYPT = RGBColor(0xC0,0x7D,0x08)
KENYA = RGBColor(0xBB,0x00,0x00)
SA = RGBColor(0x00,0x73,0x49)
CREAM = RGBColor(0xFF,0xFA,0xF0)
WARM = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xE8,0xE8,0xE8)
BLUE = RGBColor(0x19,0x76,0xD2)
GREEN = RGBColor(0x38,0x8E,0x3C)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=BLACK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=BLACK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=GOLD,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color);tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER);tb(s,1,2.8,8,0.8,sub,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);return s
def cis(fe,cn,en,color,title,items,il="📷"):
    s=ns();bg(s,CREAM);tb(s,0.3,0.15,9.4,0.6,f"{fe} {cn} {en} — {title}",sz=26,b=True,c=color,a=PP_ALIGN.CENTER)
    y=1.0
    for it,d in items:
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(y),Inches(4.8),Inches(0.4));sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
        tb(s,0.4,y+0.02,4.6,0.35,it,sz=15,b=True,c=WHITE);tb(s,0.4,y+0.45,4.7,0.5,d,sz=16,c=DARK);y+=1.1
    ib(s,5.5,1.0,4.2,y-1.1,il)
    dc={"埃及":"🏛️","肯尼亚":"🦁","南非":"🐘"}.get(cn,"🌍");tb(s,8.8,4.8,1.0,0.6,dc,sz=28,c=LGRAY,a=PP_ALIGN.RIGHT)
    return s
def vs(title,bgc):
    s=ns();bg(s,bgc);tb(s,1,0.8,8,0.8,"🎬 看视频  Watch Video",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,1.8,8,0.5,title,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
    ib(s,1.5,2.5,7,2.5,"📷 插入视频截图或粘贴视频链接");tb(s,1,5.1,8,0.3,"🔗 视频链接: ____________________",sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    return s

n=0

# 1 COVER
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.8,"Global Explorer Camp",sz=36,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.5,"环球探索沉浸式夏令营",sz=20,c=GOLD,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2),Inches(1.6),Inches(6),Inches(3.2));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GOLD;sh.line.width=Pt(3)
tf=tb(s,2.3,1.8,5.4,2.8,"BOARDING PASS  登机牌",sz=16,b=True,c=GRAY,a=PP_ALIGN.CENTER)
ap(tf,"",sz=8);ap(tf,"Flight 航班: GR EDU-002",sz=18,c=DARK);ap(tf,"Destination 目的地:  非洲 AFRICA",sz=20,b=True,c=GOLD)
ap(tf,"Date 日期: June 9, 2025",sz=16,c=DARK);ap(tf,"Gate 登机口: 谷雨大厅 GR EDU Hall",sz=16,c=DARK)
ap(tf,"",sz=8);ap(tf,"Passenger 旅客: 谷雨全体师生",sz=18,c=DARK);ap(tf,"",sz=6)
ap(tf,"Fasten seatbelts!  系好安全带！ ✈️",sz=14,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# 2 Schedule
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([("Session 1  上午","11:00-11:45","了解非洲 + 三个国家",GOLD),("Session 2  下午","2:00-2:45","复习总结 + 语言目标（认字写字）",BLUE),("Session 3  下午","3:00-4:30","写Booklet + 做Project",GREEN)]):
    y=0.9+i*1.5;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2));sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.1,4,0.4,nm,sz=20,b=True,c=WHITE);tb(s,0.7,y+0.5,3,0.4,tm,sz=16,c=RGBColor(0xFF,0xF3,0xE0));tb(s,4.5,y+0.15,4.8,0.9,dc,sz=18,c=WHITE)
pn(s,n)

# 3 Objectives
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"📚 内容目标:",sz=20,b=True,c=GOLD)
tf=tb(s,0.7,1.4,8.5,1.2,"1. 了解非洲的地理位置和特点",sz=16,c=DARK);ap(tf,"2. 了解三个国家：埃及、肯尼亚、南非（国旗、首都、景点、文化）",sz=16,c=DARK)
tb(s,0.5,2.8,9,0.5,"🗣️ 语言目标:",sz=20,b=True,c=BLUE)
tb(s,0.7,3.3,4.2,0.9,"👀 我会认：非洲 埃及 南非\n　　　　　肯尼亚 金字塔",sz=16,b=True,c=DARK)
tb(s,5.2,3.3,4.2,0.9,"✍️ 我会写：非洲 南非 金字塔",sz=16,b=True,c=DARK)
tb(s,0.5,4.3,9,0.5,"🎨 实践目标: 完成非洲Booklet + 手工项目",sz=16,c=GREEN);pn(s,n)

# 4 S1 divider
div("Session 1  上午","了解非洲的地理位置和特点\n了解三个主要国家：埃及、肯尼亚、南非",GOLD,"🌍");n+=1

# 5 Video Africa
s=vs("认识非洲 About Africa",RGBColor(0x3E,0x27,0x23));n+=1;pn(s,n)

# 6 Where is Africa
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 非洲在哪里？ Where is Africa?")
tb(s,0.4,0.9,4.5,0.5,"非洲是世界上第二大的洲！",sz=20,b=True,c=GOLD);tb(s,0.4,1.4,4.5,0.4,"Africa is the 2nd largest continent!",sz=14,c=GRAY)
ib(s,5.2,0.8,4.5,4.2,"📷 世界地图\n标出非洲位置");pn(s,n)

# 7-8 About Africa
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 认识非洲  About Africa (1/2)")
for i,(t,d) in enumerate([("🏜️ 世界第二大洲 2nd Largest","面积3,037万km² — 占地球20%"),("🌐 54个国家 54 Countries","非洲有54个国家，人口14亿！")]):
    tb(s,0.4,0.9+i*1.2,4.8,0.5,t,sz=18,b=True,c=GOLD);tb(s,0.4,1.4+i*1.2,4.8,0.4,d,sz=15,c=DARK)
ib(s,5.5,0.9,4.2,3.5,"📷 非洲地图/风景");pn(s,n)

s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 认识非洲  About Africa (2/2)")
for i,(t,d) in enumerate([("🏜️ 撒哈拉沙漠 Sahara Desert","世界上最大的热沙漠！"),("🦁 野生动物 Wildlife","狮子、大象、长颈鹿...非洲大草原！")]):
    tb(s,0.4,0.9+i*1.2,4.8,0.5,t,sz=18,b=True,c=GOLD);tb(s,0.4,1.4+i*1.2,4.8,0.4,d,sz=15,c=DARK)
ib(s,5.5,0.9,4.2,3.5,"📷 撒哈拉沙漠/非洲动物");pn(s,n)

# EGYPT 埃及
s=vs("🇪🇬 认识埃及 About Egypt",RGBColor(0x5D,0x40,0x00));n+=1;pn(s,n)
s=cis("🇪🇬","埃及","Egypt",EGYPT,"国旗 + 首都",[("🏴 国旗","红白黑三色+金色鹰 Red/White/Black+Eagle"),("🏛️ 首都","开罗 Cairo")],"📷 埃及国旗+开罗");n+=1;pn(s,n)
s=cis("🇪🇬","埃及","Egypt",EGYPT,"人口 + 语言",[("👥 人口","约1.1亿 ~110 million"),("🗣️ 语言","阿拉伯语 Arabic")],"📷 埃及图片");n+=1;pn(s,n)
s=cis("🇪🇬","埃及","Egypt",EGYPT,"主要景点 Landmarks",[("🏛️ 金字塔 Pyramids","世界七大奇迹之一！4,500年历史"),("🗿 狮身人面像 Sphinx","守护金字塔的神秘雕像"),("🏞️ 尼罗河 Nile","世界最长河流！6,650km")],"📷 金字塔/尼罗河");n+=1;pn(s,n)
s=cis("🇪🇬","埃及","Egypt",EGYPT,"礼节 Etiquette",[("👋 说你好","「Marhaba」(mar-ha-ba)"),("🤝 打招呼","握手+右手放胸口表示尊重"),("🍽️ 吃饭礼节","用右手吃饭/主人会不停添饭=热情！")],"📷 埃及问候/用餐");n+=1;pn(s,n)
s=cis("🇪🇬","埃及","Egypt",EGYPT,"美食 Food",[("🫓 大饼 Pita Bread","埃及人每天都吃的面饼！"),("🥙 Falafel 法拉费尔","炸鹰嘴豆丸子"),("🍵 薄荷茶 Mint Tea","又甜又香的国民饮料")],"📷 埃及美食");n+=1;pn(s,n)

# KENYA 肯尼亚
s=vs("🇰🇪 认识肯尼亚 About Kenya",RGBColor(0x5D,0x00,0x00));n+=1;pn(s,n)
s=cis("🇰🇪","肯尼亚","Kenya",KENYA,"国旗 + 首都",[("🏴 国旗","黑红绿三色+马赛盾牌 Black/Red/Green+Shield"),("🏛️ 首都","内罗毕 Nairobi")],"📷 肯尼亚国旗");n+=1;pn(s,n)
s=cis("🇰🇪","肯尼亚","Kenya",KENYA,"人口 + 语言",[("👥 人口","约5,400万 ~54 million"),("🗣️ 语言","斯瓦希里语 Swahili + 英语 English")],"📷 肯尼亚图片");n+=1;pn(s,n)
s=cis("🇰🇪","肯尼亚","Kenya",KENYA,"主要景点 Landmarks",[("🦁 马赛马拉 Maasai Mara","世界最著名野生动物保护区！"),("🏔️ 肯尼亚山 Mt. Kenya","非洲第二高峰 5,199米"),("🦒 动物大迁徙 Great Migration","每年200万动物迁徙！")],"📷 大草原/迁徙");n+=1;pn(s,n)
s=cis("🇰🇪","肯尼亚","Kenya",KENYA,"礼节 Etiquette",[("👋 说你好","「Jambo」(jam-bo) 斯瓦希里语"),("🤝 打招呼","握手 Handshake（时间比较长）"),("🍽️ 吃饭礼节","用右手/客人先吃/Ugali用手捏着吃")],"📷 肯尼亚问候");n+=1;pn(s,n)
s=cis("🇰🇪","肯尼亚","Kenya",KENYA,"美食 Food",[("🫓 Ugali 乌伽利","玉米面主食，像糍粑"),("🥩 Nyama Choma 烤肉","肯尼亚最受欢迎的美食！"),("☕ 肯尼亚咖啡","世界顶级咖啡产地之一")],"📷 肯尼亚美食");n+=1;pn(s,n)

# SOUTH AFRICA 南非
s=vs("🇿🇦 认识南非 About South Africa",RGBColor(0x00,0x3D,0x25));n+=1;pn(s,n)
s=cis("🇿🇦","南非","South Africa",SA,"国旗 + 首都",[("🏴 国旗","六种颜色！彩虹之国 Rainbow Nation 🌈"),("🏛️ 首都","比勒陀利亚 Pretoria")],"📷 南非国旗");n+=1;pn(s,n)
s=cis("🇿🇦","南非","South Africa",SA,"人口 + 语言",[("👥 人口","约6,000万 ~60 million"),("🗣️ 语言","11种官方语言！英语+祖鲁语Zulu+...")],"📷 南非图片");n+=1;pn(s,n)
s=cis("🇿🇦","南非","South Africa",SA,"主要景点 Landmarks",[("🏔️ 桌山 Table Mountain","开普敦标志，平坦的山顶！"),("🐧 企鹅海滩 Boulders Beach","非洲也有企鹅！"),("🦏 克鲁格公园 Kruger Park","非洲最大野生动物保护区之一")],"📷 桌山/企鹅/克鲁格");n+=1;pn(s,n)
s=cis("🇿🇦","南非","South Africa",SA,"礼节 Etiquette",[("👋 说你好","「Sawubona」(sa-wu-bo-na) 祖鲁语"),("🤝 打招呼","握手+有时碰拳 Fist bump"),("🍽️ 吃饭礼节","Braai烧烤是社交活动/用手吃传统食物")],"📷 南非问候/Braai");n+=1;pn(s,n)
s=cis("🇿🇦","南非","South Africa",SA,"美食 Food",[("🥩 Braai 烧烤","南非国民活动！类似BBQ"),("🫓 Bobotie 波波提","咖喱肉末砂锅，南非国菜"),("🍬 Biltong 牛肉干","南非特色零食")],"📷 南非美食");n+=1;pn(s,n)

# Greeting Practice
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 打招呼练习  Greeting Practice")
for i,(nm,gr,cl) in enumerate([("🇪🇬 埃及","握手+「Marhaba」\n(mar-ha-ba)",EGYPT),("🇰🇪 肯尼亚","握手+「Jambo」\n(jam-bo)",KENYA),("🇿🇦 南非","握手/碰拳+「Sawubona」\n(sa-wu-bo-na)",SA)]):
    x=0.4+i*3.2;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.9),Inches(3.5));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.1,2.7,0.5,nm,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER);ls=gr.split('\n');tf=tb(s,x+0.1,1.7,2.7,0.3,ls[0],sz=14,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=14,c=DARK,a=PP_ALIGN.CENTER)
    ib(s,x+0.3,2.4,2.3,1.8,"📷 动作示范")
tb(s,0.4,4.7,9,0.4,"站起来和旁边的同学练习！Stand up and practice!",sz=14,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# SESSION 2
div("Session 2  下午","复习总结 + 语言目标\n我会认：非洲 埃及 南非 肯尼亚 金字塔\n我会写：非洲 南非 金字塔",BLUE,"📖");n+=1

# Comparison blank
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 非洲三国文化对比  你知道吗？",BLUE)
tb(s,0.4,0.85,9,0.35,"你能填出来吗？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(7,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.5);t.columns[2].width=Inches(2.5);t.columns[3].width=Inches(2.6)
for r,rd in enumerate([["","🇪🇬 埃及","🇰🇪 肯尼亚","🇿🇦 南非"],["👋 说你好","？","？","？"],["🤝 打招呼","？","？","？"],["🍽️ 美食","？","？","？"],["🏛️ 首都","？","？","？"],["🏔️ 景点","？","？","？"],["🐾 动物","？","？","？"]]):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        if r==0:rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=BLUE
        elif c==0:rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        else:rn.font.color.rgb=LGRAY;rn.font.size=Pt(20)
pn(s,n)

# Comparison answers
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 非洲三国文化对比  Comparison",BLUE)
tb(s,0.4,0.85,9,0.35,"三个国家各有特色，你最想去哪个？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(7,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.5);t.columns[2].width=Inches(2.5);t.columns[3].width=Inches(2.6)
for r,rd in enumerate([["","🇪🇬 埃及","🇰🇪 肯尼亚","🇿🇦 南非"],
    ["👋 说你好","Marhaba\n(mar-ha-ba)","Jambo\n(jam-bo)","Sawubona\n(sa-wu-bo-na)"],
    ["🤝 打招呼","握手+手放胸口","握手(较长)","握手/碰拳"],
    ["🍽️ 美食","大饼Pita\nFalafel","Ugali乌伽利\n烤肉","Braai烧烤\nBobotie"],
    ["🏛️ 首都","开罗 Cairo","内罗毕 Nairobi","比勒陀利亚 Pretoria"],
    ["🏔️ 景点","金字塔 Pyramids","马赛马拉 Mara","桌山 Table Mt."],
    ["🐾 动物","🐪 骆驼","🦁 狮子","🐘 大象"]]):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        rn.font.color.rgb=WHITE if r==0 else(DARK if c>0 else RGBColor(0x44,0x44,0x44))
        if r==0:cl.fill.solid();cl.fill.fore_color.rgb=BLUE
        elif c==0:cl.fill.solid();cl.fill.fore_color.rgb=WARM
        elif r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF5)
pn(s,n)

# Word cards 我会认
for w,py,en,sent,il in [("非洲","fēi zhōu","Africa","非洲是世界上第二大的洲。","📷 非洲地图"),("埃及","āi jí","Egypt","埃及有著名的金字塔。","📷 金字塔"),("南非","nán fēi","South Africa","南非被称为彩虹之国。","📷 南非桌山"),("肯尼亚","kěn ní yà","Kenya","肯尼亚有很多野生动物。","📷 大草原"),("金字塔","jīn zì tǎ","Pyramid","金字塔有四千五百年的历史。","📷 金字塔")]:
    s=ns();n+=1;bg(s,CREAM);hb(s,"👀 我会认  I Can Read",BLUE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5));sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=GOLD,a=PP_ALIGN.CENTER);tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=BLUE,a=PP_ALIGN.CENTER);ib(s,5.3,1.0,4.4,2.5,il)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2));sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=BLUE;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=BLUE);tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK);pn(s,n)

# Word games
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",BLUE)
for i,(nm,name,desc,bgc) in enumerate([("1️⃣","拍苍蝇 Fly Swatter","把字卡贴在白板上\n老师说词语，学生拍！",WARM),("2️⃣","举牌游戏 Show Me","每人5张字卡\n老师说词语，举正确的卡",RGBColor(0xE3,0xF2,0xFD)),("3️⃣","抢椅子 Musical Chairs","椅子上放字卡\n音乐停，读出词",RGBColor(0xE8,0xF5,0xE9)),("4️⃣","传话筒 Pass the Mic","传球，停下的人\n读字卡并造句",RGBColor(0xFC,0xE4,0xEC))]):
    x=0.3+i*2.4;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2));sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,nm,sz=24,a=PP_ALIGN.CENTER);tb(s,x+0.1,1.4,2.0,0.6,name,sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls=desc.split('\n');tf=tb(s,x+0.15,2.1,1.9,1.5,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.8,2.0,0.4,"低prep ✅",sz=12,b=True,c=GREEN,a=PP_ALIGN.CENTER)
pn(s,n)

# Write cards 我会写
for w,py,en,il in [("非洲","fēi zhōu","Africa","📷 非洲地图"),("南非","nán fēi","South Africa","📷 南非国旗"),("金字塔","jīn zì tǎ","Pyramid","📷 金字塔")]:
    s=ns();n+=1;bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",BLUE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BLUE;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=BLUE,a=PP_ALIGN.CENTER);tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,il)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8));sh2.fill.solid();sh2.fill.fore_color.rgb=RGBColor(0xE3,0xF2,0xFD);sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=BLUE);ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=BLUE);ap(tf,"1. 空中写 Air Write",sz=13,c=DARK);ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK);ap(tf,"3. 纸上写 Write 3 times",sz=13,c=DARK)
    pn(s,n)

# SESSION 3
div("Session 3  下午","写Booklet + 做Project\n3:00 - 4:30",GREEN,"🎨");n+=1

# Booklet
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成"探索非洲"练习册',GREEN);ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图");pn(s,n)

# Project overview
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 Project Time!  4个手工项目",GREEN)
for i,(pr,nm,md,bc) in enumerate([("PROJECT 1","🧩 非洲拼图","group project",WARM),("PROJECT 2","🦁 动物面具","独立完成",RGBColor(0xE3,0xF2,0xFD)),("PROJECT 3","🏛️ 建金字塔","group project",RGBColor(0xE8,0xF5,0xE9)),("PROJECT 4","📿 文化手环","独立完成",RGBColor(0xFC,0xE4,0xEC))]):
    x=0.3+i*2.4;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2));sh.fill.solid();sh.fill.fore_color.rgb=bc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.3,pr,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER);tb(s,x+0.1,1.3,2.0,0.6,nm,sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.9,2.0,0.3,f"({md})",sz=11,c=GRAY,a=PP_ALIGN.CENTER);ib(s,x+0.2,2.4,1.8,2.3,"📷 示范")
pn(s,n)

# Individual projects
for pr,nm,md,bc,cl in [("PROJECT 1","🧩 非洲拼图  Africa Puzzle Map","group project",WARM,GREEN),("PROJECT 2","🦁 非洲动物面具  Animal Mask","独立完成",RGBColor(0xE3,0xF2,0xFD),BLUE),("PROJECT 3","🏛️ 建金字塔  Build a Pyramid","group project",RGBColor(0xE8,0xF5,0xE9),GREEN),("PROJECT 4","📿 非洲文化手环  Pattern Bracelet","独立完成",RGBColor(0xFC,0xE4,0xEC),RGBColor(0xC2,0x18,0x5B))]:
    s=ns();n+=1;bg(s,bc);hb(s,f"{pr}: {nm}",cl);tb(s,0.4,0.85,9,0.3,f"({md})",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,0.4,1.3,4.4,3.5,"📷 示范图片/截图");ib(s,5.2,1.3,4.5,3.5,"🎬 教学视频/步骤视频");pn(s,n)

# Visa stamp
s=ns();n+=1;bg(s,CREAM);tb(s,1,0.5,8,0.8,"🪪 非洲签证章  Africa Visa Stamp",sz=30,b=True,c=GOLD,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.5),Inches(3),Inches(3));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GOLD;sh.line.width=Pt(5)
tf=tb(s,3.6,1.8,2.8,2.5,"AFRICA\n非洲",sz=28,b=True,c=GOLD,a=PP_ALIGN.CENTER);ap(tf,"✓ VISITED",sz=16,b=True,c=GREEN,a=PP_ALIGN.CENTER);ap(tf,"6/9/2025",sz=12,c=GRAY,a=PP_ALIGN.CENTER);ap(tf,"埃及 · 肯尼亚 · 南非",sz=12,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1,4.7,8,0.4,"恭喜你完成非洲之旅！Congratulations! 🎉",sz=16,b=True,c=GOLD,a=PP_ALIGN.CENTER);pn(s,n)

# Tomorrow
s=ns();n+=1;bg(s,GOLD)
tb(s,1,1.0,8,0.8,"✈️ 明天航班  Tomorrow's Flight",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,2,2.2,6,2.5,"Flight 航班: GR EDU-003",sz=22,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Destination 目的地: 欧洲 EUROPE 🌍",sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER);ap(tf,"",sz=10)
ap(tf,"明天我们去欧洲！",sz=20,c=WHITE,a=PP_ALIGN.CENTER);ap(tf,"那里有什么有名的地方？",sz=18,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
ap(tf,"",sz=10);ap(tf,"See you tomorrow, explorers! 明天见！",sz=16,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/世界旅行world_trip_pbl/day2_africa_v2.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
