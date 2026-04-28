#!/usr/bin/env python3
"""
Day 3 Europe PPT — Same structure as Day 2 Africa.
Countries: France 法国, Italy 意大利, UK 英国
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# Colors
EUROPE_BLUE = RGBColor(0x1A,0x23,0x7E)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BLACK = RGBColor(0x33,0x33,0x33)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
FRANCE = RGBColor(0x00,0x23,0x95)
ITALY = RGBColor(0x00,0x8C,0x45)
UK = RGBColor(0xC8,0x10,0x2E)
CREAM = RGBColor(0xFF,0xFA,0xF0)
WARM = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xE8,0xE8,0xE8)
BLUE = RGBColor(0x19,0x76,0xD2)
GREEN = RGBColor(0x38,0x8E,0x3C)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=BLACK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=BLACK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=EUROPE_BLUE,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color);tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER);tb(s,1,2.8,8,0.8,sub,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);return s
def cis(fe,cn,en,color,title,items,il="📷"):
    s=ns();bg(s,CREAM);tb(s,0.3,0.15,9.4,0.6,f"{fe} {cn} {en} — {title}",sz=26,b=True,c=color,a=PP_ALIGN.CENTER)
    y=1.0
    for it,d in items:
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(y),Inches(4.8),Inches(0.4));sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
        tb(s,0.4,y+0.02,4.6,0.35,it,sz=15,b=True,c=WHITE);tb(s,0.4,y+0.45,4.7,0.5,d,sz=16,c=DARK);y+=1.1
    ib(s,5.5,1.0,4.2,y-1.1,il)
    dc={"法国":"🗼","意大利":"🍕","英国":"🎡"}.get(cn,"🌍");tb(s,8.8,4.8,1.0,0.6,dc,sz=28,c=LGRAY,a=PP_ALIGN.RIGHT)
    return s
def vs(title,bgc):
    s=ns();bg(s,bgc);tb(s,1,0.8,8,0.8,"🎬 看视频  Watch Video",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,1.8,8,0.5,title,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
    ib(s,1.5,2.5,7,2.5,"📷 插入视频截图或粘贴视频链接");tb(s,1,5.1,8,0.3,"🔗 视频链接: ____________________",sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    return s

n=0

# 1 COVER
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.8,"Global Explorer Camp",sz=36,b=True,c=EUROPE_BLUE,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.5,"环球探索沉浸式夏令营",sz=20,c=EUROPE_BLUE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2),Inches(1.6),Inches(6),Inches(3.2));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=EUROPE_BLUE;sh.line.width=Pt(3)
tf=tb(s,2.3,1.8,5.4,2.8,"BOARDING PASS  登机牌",sz=16,b=True,c=GRAY,a=PP_ALIGN.CENTER)
ap(tf,"",sz=8);ap(tf,"Flight 航班: GR EDU-003",sz=18,c=DARK);ap(tf,"Destination 目的地:  欧洲 EUROPE",sz=20,b=True,c=EUROPE_BLUE)
ap(tf,"Date 日期: June 10, 2025",sz=16,c=DARK);ap(tf,"Gate 登机口: 谷雨大厅 GR EDU Hall",sz=16,c=DARK)
ap(tf,"",sz=8);ap(tf,"Passenger 旅客: 谷雨全体师生",sz=18,c=DARK);ap(tf,"",sz=6)
ap(tf,"Fasten seatbelts!  系好安全带！ ✈️",sz=14,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# 2 Schedule
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([("Session 1  上午","11:00-11:45","了解欧洲 + 三个国家",EUROPE_BLUE),("Session 2  下午","2:00-2:45","复习总结 + 语言目标（认字写字）",BLUE),("Session 3  下午","3:00-4:30","写Booklet + 做Project",GREEN)]):
    y=0.9+i*1.5;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2));sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.1,4,0.4,nm,sz=20,b=True,c=WHITE);tb(s,0.7,y+0.5,3,0.4,tm,sz=16,c=RGBColor(0xFF,0xF3,0xE0));tb(s,4.5,y+0.15,4.8,0.9,dc,sz=18,c=WHITE)
pn(s,n)

# 3 Objectives
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"📚 内容目标:",sz=20,b=True,c=EUROPE_BLUE)
tf=tb(s,0.7,1.4,8.5,1.2,"1. 了解欧洲的地理位置和特点",sz=16,c=DARK);ap(tf,"2. 了解三个国家：法国、意大利、英国（国旗、首都、景点、文化）",sz=16,c=DARK)
tb(s,0.5,2.8,9,0.5,"🗣️ 语言目标:",sz=20,b=True,c=BLUE)
tb(s,0.7,3.3,4.2,0.9,"👀 我会认：欧洲 法国 意大利 英国",sz=16,b=True,c=DARK)
tb(s,5.2,3.3,4.2,0.9,"✍️ 我会写：欧洲 法国 英国",sz=16,b=True,c=DARK)
tb(s,0.5,4.3,9,0.5,"🎨 实践目标: 完成欧洲Booklet + 手工项目",sz=16,c=GREEN);pn(s,n)

# 4 S1 divider
div("Session 1  上午","了解欧洲的地理位置和特点\n了解三个主要国家：法国、意大利、英国",EUROPE_BLUE,"🌍");n+=1

# 5 Video Europe
s=vs("认识欧洲 About Europe",RGBColor(0x0D,0x14,0x52));n+=1;pn(s,n)

# 6 Where is Europe
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 欧洲在哪里？ Where is Europe?")
tb(s,0.4,0.9,4.5,0.5,"欧洲是第二小的洲但历史最悠久之一！",sz=20,b=True,c=EUROPE_BLUE);tb(s,0.4,1.4,4.5,0.4,"Europe is the 2nd smallest but one of the oldest!",sz=14,c=GRAY)
ib(s,5.2,0.8,4.5,4.2,"📷 世界地图\n标出欧洲位置");pn(s,n)

# 7-8 About Europe
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 认识欧洲  About Europe (1/2)")
for i,(t,d) in enumerate([("🏰 44个国家 44 Countries","欧洲有44个国家，人口7.5亿！"),("🌐 历史悠久 Rich History","文艺复兴、工业革命的发源地")]):
    tb(s,0.4,0.9+i*1.2,4.8,0.5,t,sz=18,b=True,c=EUROPE_BLUE);tb(s,0.4,1.4+i*1.2,4.8,0.4,d,sz=15,c=DARK)
ib(s,5.5,0.9,4.2,3.5,"📷 欧洲地图/风景");pn(s,n)

s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 认识欧洲  About Europe (2/2)")
for i,(t,d) in enumerate([("🏔️ 阿尔卑斯山 Alps","欧洲最高的山脉，横跨八个国家！"),("🌊 多瑙河 Danube River","欧洲第二长河，流经十个国家！")]):
    tb(s,0.4,0.9+i*1.2,4.8,0.5,t,sz=18,b=True,c=EUROPE_BLUE);tb(s,0.4,1.4+i*1.2,4.8,0.4,d,sz=15,c=DARK)
ib(s,5.5,0.9,4.2,3.5,"📷 阿尔卑斯山/多瑙河");pn(s,n)

# FRANCE 法国
s=vs("🇫🇷 认识法国 About France",RGBColor(0x00,0x14,0x5A));n+=1;pn(s,n)
s=cis("🇫🇷","法国","France",FRANCE,"国旗 + 首都",[("🏴 国旗","蓝白红三色 Blue/White/Red"),("🏛️ 首都","巴黎 Paris")],"📷 法国国旗+巴黎");n+=1;pn(s,n)
s=cis("🇫🇷","法国","France",FRANCE,"人口 + 语言",[("👥 人口","约6,700万 ~67 million"),("🗣️ 语言","法语 French")],"📷 法国图片");n+=1;pn(s,n)
s=cis("🇫🇷","法国","France",FRANCE,"主要景点 Landmarks",[("🗼 埃菲尔铁塔 Eiffel Tower","巴黎地标，建于1889年！"),("🖼️ 卢浮宫 Louvre","世界最大的艺术博物馆"),("🏰 凡尔赛宫 Versailles","法国国王的奢华宫殿")],"📷 埃菲尔铁塔/卢浮宫");n+=1;pn(s,n)
s=cis("🇫🇷","法国","France",FRANCE,"礼节 Etiquette",[("👋 说你好","「Bonjour」(bon-zhur)"),("🤝 打招呼","贴面礼 Cheek kiss（左右各一次）"),("🍽️ 吃饭礼节","面包配每餐/不要急着吃=慢慢享受")],"📷 法国问候/用餐");n+=1;pn(s,n)
s=cis("🇫🇷","法国","France",FRANCE,"美食 Food",[("🥐 可颂 Croissant","法国经典早餐面包！"),("🥖 法棍 Baguette","又长又脆的法式面包"),("🧁 马卡龙 Macaron","五颜六色的法式甜点")],"📷 法国美食");n+=1;pn(s,n)

# ITALY 意大利
s=vs("🇮🇹 认识意大利 About Italy",RGBColor(0x00,0x4D,0x25));n+=1;pn(s,n)
s=cis("🇮🇹","意大利","Italy",ITALY,"国旗 + 首都",[("🏴 国旗","绿白红三色 Green/White/Red"),("🏛️ 首都","罗马 Rome")],"📷 意大利国旗+罗马");n+=1;pn(s,n)
s=cis("🇮🇹","意大利","Italy",ITALY,"人口 + 语言",[("👥 人口","约5,900万 ~59 million"),("🗣️ 语言","意大利语 Italian")],"📷 意大利图片");n+=1;pn(s,n)
s=cis("🇮🇹","意大利","Italy",ITALY,"主要景点 Landmarks",[("🏗️ 比萨斜塔 Leaning Tower of Pisa","世界著名的倾斜建筑！"),("🏟️ 罗马斗兽场 Colosseum","古罗马最大的竞技场"),("🚣 威尼斯 Venice水城","用船代替汽车的城市！")],"📷 比萨斜塔/斗兽场");n+=1;pn(s,n)
s=cis("🇮🇹","意大利","Italy",ITALY,"礼节 Etiquette",[("👋 说你好","「Ciao」(chao)"),("🤝 打招呼","拥抱+贴面 Hug+cheek kiss"),("🍽️ 吃饭礼节","意大利面要用叉子卷着吃/不要加ketchup!")],"📷 意大利问候/用餐");n+=1;pn(s,n)
s=cis("🇮🇹","意大利","Italy",ITALY,"美食 Food",[("🍕 披萨 Pizza","意大利最著名的美食！"),("🍝 意大利面 Pasta","各种形状和酱料"),("🍦 冰淇淋 Gelato","比普通冰淇淋更浓郁！")],"📷 意大利美食");n+=1;pn(s,n)

# UK 英国
s=vs("🇬🇧 认识英国 About UK",RGBColor(0x5D,0x08,0x15));n+=1;pn(s,n)
s=cis("🇬🇧","英国","UK",UK,"国旗 + 首都",[("🏴 国旗","红白蓝米字旗 Union Jack"),("🏛️ 首都","伦敦 London")],"📷 英国国旗+伦敦");n+=1;pn(s,n)
s=cis("🇬🇧","英国","UK",UK,"人口 + 语言",[("👥 人口","约6,700万 ~67 million"),("🗣️ 语言","英语 English")],"📷 英国图片");n+=1;pn(s,n)
s=cis("🇬🇧","英国","UK",UK,"主要景点 Landmarks",[("🕐 大本钟 Big Ben","伦敦最著名的钟楼！"),("🏰 白金汉宫 Buckingham Palace","英国女王的官邸"),("🎡 伦敦眼 London Eye","泰晤士河畔的巨型摩天轮")],"📷 大本钟/白金汉宫");n+=1;pn(s,n)
s=cis("🇬🇧","英国","UK",UK,"礼节 Etiquette",[("👋 说你好","「Hello」"),("🤝 打招呼","握手 Handshake（保持距离）"),("🍽️ 吃饭礼节","下午茶文化 Afternoon tea/刀叉用法讲究")],"📷 英国问候/下午茶");n+=1;pn(s,n)
s=cis("🇬🇧","英国","UK",UK,"美食 Food",[("🐟 炸鱼薯条 Fish & Chips","英国最经典的街头美食！"),("🍳 英式早餐 Full English Breakfast","鸡蛋+培根+豆子+香肠"),("🫖 司康饼 Scone","配下午茶的经典点心")],"📷 英国美食");n+=1;pn(s,n)

# Greeting Practice
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 打招呼练习  Greeting Practice")
for i,(nm,gr,cl) in enumerate([("🇫🇷 法国","贴面礼+「Bonjour」\n(bon-zhur)",FRANCE),("🇮🇹 意大利","拥抱+贴面+「Ciao」\n(chao)",ITALY),("🇬🇧 英国","握手+「Hello」\n(保持距离)",UK)]):
    x=0.4+i*3.2;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.9),Inches(3.5));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.1,2.7,0.5,nm,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER);ls=gr.split('\n');tf=tb(s,x+0.1,1.7,2.7,0.3,ls[0],sz=14,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=14,c=DARK,a=PP_ALIGN.CENTER)
    ib(s,x+0.3,2.4,2.3,1.8,"📷 动作示范")
tb(s,0.4,4.7,9,0.4,"站起来和旁边的同学练习！Stand up and practice!",sz=14,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# SESSION 2
div("Session 2  下午","复习总结 + 语言目标\n我会认：欧洲 法国 意大利 英国\n我会写：欧洲 法国 英国",BLUE,"📖");n+=1

# Comparison blank
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 欧洲三国文化对比  你知道吗？",BLUE)
tb(s,0.4,0.85,9,0.35,"你能填出来吗？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(7,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.5);t.columns[2].width=Inches(2.5);t.columns[3].width=Inches(2.6)
for r,rd in enumerate([["","🇫🇷 法国","🇮🇹 意大利","🇬🇧 英国"],["👋 说你好","？","？","？"],["🤝 打招呼","？","？","？"],["🍽️ 美食","？","？","？"],["🏛️ 首都","？","？","？"],["🏔️ 景点","？","？","？"],["🎁 代表物","？","？","？"]]):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        if r==0:rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=BLUE
        elif c==0:rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        else:rn.font.color.rgb=LGRAY;rn.font.size=Pt(20)
pn(s,n)

# Comparison answers
s=ns();n+=1;bg(s,CREAM);hb(s,"🌍 欧洲三国文化对比  Comparison",BLUE)
tb(s,0.4,0.85,9,0.35,"三个国家各有特色，你最想去哪个？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(7,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.5);t.columns[2].width=Inches(2.5);t.columns[3].width=Inches(2.6)
for r,rd in enumerate([["","🇫🇷 法国","🇮🇹 意大利","🇬🇧 英国"],
    ["👋 说你好","Bonjour\n(bon-zhur)","Ciao\n(chao)","Hello"],
    ["🤝 打招呼","贴面礼\n(左右各一次)","拥抱+贴面","握手(保持距离)"],
    ["🍽️ 美食","可颂Croissant\n马卡龙Macaron","披萨Pizza\nGelato冰淇淋","Fish & Chips\nScone司康饼"],
    ["🏛️ 首都","巴黎 Paris","罗马 Rome","伦敦 London"],
    ["🏔️ 景点","埃菲尔铁塔\nEiffel Tower","斗兽场\nColosseum","大本钟\nBig Ben"],
    ["🎁 代表物","🗼 铁塔","🍕 披萨","🎡 伦敦眼"]]):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        rn.font.color.rgb=WHITE if r==0 else(DARK if c>0 else RGBColor(0x44,0x44,0x44))
        if r==0:cl.fill.solid();cl.fill.fore_color.rgb=BLUE
        elif c==0:cl.fill.solid();cl.fill.fore_color.rgb=WARM
        elif r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF5)
pn(s,n)

# Word cards 我会认
for w,py,en,sent,il in [("欧洲","ōu zhōu","Europe","欧洲有很多古老的城堡。","📷 欧洲地图"),("法国","fǎ guó","France","法国的首都是巴黎。","📷 埃菲尔铁塔"),("意大利","yì dà lì","Italy","意大利的披萨很有名。","📷 比萨斜塔"),("英国","yīng guó","UK","英国人喜欢喝下午茶。","📷 大本钟")]:
    s=ns();n+=1;bg(s,CREAM);hb(s,"👀 我会认  I Can Read",BLUE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5));sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=EUROPE_BLUE,a=PP_ALIGN.CENTER);tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=BLUE,a=PP_ALIGN.CENTER);ib(s,5.3,1.0,4.4,2.5,il)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2));sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=BLUE;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=BLUE);tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK);pn(s,n)

# Word games
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",BLUE)
for i,(nm,name,desc,bgc) in enumerate([("1️⃣","拍苍蝇 Fly Swatter","把字卡贴在白板上\n老师说词语，学生拍！",WARM),("2️⃣","举牌游戏 Show Me","每人4张字卡\n老师说词语，举正确的卡",RGBColor(0xE3,0xF2,0xFD)),("3️⃣","抢椅子 Musical Chairs","椅子上放字卡\n音乐停，读出词",RGBColor(0xE8,0xF5,0xE9)),("4️⃣","传话筒 Pass the Mic","传球，停下的人\n读字卡并造句",RGBColor(0xFC,0xE4,0xEC))]):
    x=0.3+i*2.4;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2));sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,nm,sz=24,a=PP_ALIGN.CENTER);tb(s,x+0.1,1.4,2.0,0.6,name,sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls=desc.split('\n');tf=tb(s,x+0.15,2.1,1.9,1.5,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.8,2.0,0.4,"低prep ✅",sz=12,b=True,c=GREEN,a=PP_ALIGN.CENTER)
pn(s,n)

# Write cards 我会写
for w,py,en,il in [("欧洲","ōu zhōu","Europe","📷 欧洲地图"),("法国","fǎ guó","France","📷 法国国旗"),("英国","yīng guó","UK","📷 英国国旗")]:
    s=ns();n+=1;bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",BLUE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BLUE;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=BLUE,a=PP_ALIGN.CENTER);tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,il)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8));sh2.fill.solid();sh2.fill.fore_color.rgb=RGBColor(0xE3,0xF2,0xFD);sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=BLUE);ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=BLUE);ap(tf,"1. 空中写 Air Write",sz=13,c=DARK);ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK);ap(tf,"3. 纸上写 Write 3 times",sz=13,c=DARK)
    pn(s,n)

# SESSION 3
div("Session 3  下午","写Booklet + 做Project\n3:00 - 4:30",GREEN,"🎨");n+=1

# Booklet
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成"探索欧洲"练习册',GREEN);ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图");pn(s,n)

# Project overview
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 Project Time!  4个手工项目",GREEN)
for i,(pr,nm,md,bc) in enumerate([("PROJECT 1","🧩 欧洲拼图","group project",WARM),("PROJECT 2","🏛️ 地标手工","独立完成",RGBColor(0xE3,0xF2,0xFD)),("PROJECT 3","🍕 欧洲美食艺术","独立完成",RGBColor(0xE8,0xF5,0xE9)),("PROJECT 4","🎨 意大利面艺术","group project",RGBColor(0xFC,0xE4,0xEC))]):
    x=0.3+i*2.4;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2));sh.fill.solid();sh.fill.fore_color.rgb=bc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.3,pr,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER);tb(s,x+0.1,1.3,2.0,0.6,nm,sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.9,2.0,0.3,f"({md})",sz=11,c=GRAY,a=PP_ALIGN.CENTER);ib(s,x+0.2,2.4,1.8,2.3,"📷 示范")
pn(s,n)

# Individual projects
for pr,nm,md,bc,cl in [("PROJECT 1","🧩 欧洲拼图  Europe Puzzle Map","group project",WARM,GREEN),("PROJECT 2","🏛️ 地标手工  Landmark Craft","独立完成",RGBColor(0xE3,0xF2,0xFD),BLUE),("PROJECT 3","🍕 欧洲美食艺术  Food Craft","独立完成",RGBColor(0xE8,0xF5,0xE9),GREEN),("PROJECT 4","🎨 意大利面艺术  Pasta Art","group project",RGBColor(0xFC,0xE4,0xEC),RGBColor(0xC2,0x18,0x5B))]:
    s=ns();n+=1;bg(s,bc);hb(s,f"{pr}: {nm}",cl);tb(s,0.4,0.85,9,0.3,f"({md})",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,0.4,1.3,4.4,3.5,"📷 示范图片/截图");ib(s,5.2,1.3,4.5,3.5,"🎬 教学视频/步骤视频");pn(s,n)

# Visa stamp
s=ns();n+=1;bg(s,CREAM);tb(s,1,0.5,8,0.8,"🪪 欧洲签证章  Europe Visa Stamp",sz=30,b=True,c=EUROPE_BLUE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.5),Inches(3),Inches(3));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=EUROPE_BLUE;sh.line.width=Pt(5)
tf=tb(s,3.6,1.8,2.8,2.5,"EUROPE\n欧洲",sz=28,b=True,c=EUROPE_BLUE,a=PP_ALIGN.CENTER);ap(tf,"✓ VISITED",sz=16,b=True,c=GREEN,a=PP_ALIGN.CENTER);ap(tf,"6/10/2025",sz=12,c=GRAY,a=PP_ALIGN.CENTER);ap(tf,"法国 · 意大利 · 英国",sz=12,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1,4.7,8,0.4,"恭喜你完成欧洲之旅！Congratulations! 🎉",sz=16,b=True,c=EUROPE_BLUE,a=PP_ALIGN.CENTER);pn(s,n)

# Tomorrow
s=ns();n+=1;bg(s,EUROPE_BLUE)
tb(s,1,1.0,8,0.8,"✈️ 明天航班  Tomorrow's Flight",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,2,2.2,6,2.5,"Flight 航班: GR EDU-004",sz=22,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Destination 目的地: 美洲 AMERICAS 🌎",sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER);ap(tf,"",sz=10)
ap(tf,"明天我们去美洲！",sz=20,c=WHITE,a=PP_ALIGN.CENTER);ap(tf,"那里有什么有名的地方？",sz=18,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
ap(tf,"",sz=10);ap(tf,"See you tomorrow, explorers! 明天见！",sz=16,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/世界旅行world_trip_pbl/day3_europe_v2.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
