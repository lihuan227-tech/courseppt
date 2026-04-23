#!/usr/bin/env python3
"""
Day 4 Americas PPT — Same structure as Day 2 Africa v2.
Countries: USA 美国, Canada 加拿大, Mexico 墨西哥
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# Colors
TEAL = RGBColor(0x00,0x69,0x5C)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BLACK = RGBColor(0x33,0x33,0x33)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
USA = RGBColor(0x3C,0x3B,0x6E)
CANADA = RGBColor(0xD5,0x21,0x1E)
MEXICO = RGBColor(0x00,0x6D,0x47)
CREAM = RGBColor(0xFF,0xFA,0xF0)
WARM = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xE8,0xE8,0xE8)
BLUE = RGBColor(0x19,0x76,0xD2)
GREEN = RGBColor(0x38,0x8E,0x3C)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=BLACK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=BLACK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=TEAL,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color);tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER);tb(s,1,2.8,8,0.8,sub,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);return s
def cis(fe,cn,en,color,title,items,il="📷"):
    s=ns();bg(s,CREAM);tb(s,0.3,0.15,9.4,0.6,f"{fe} {cn} {en} — {title}",sz=26,b=True,c=color,a=PP_ALIGN.CENTER)
    y=1.0
    for it,d in items:
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(y),Inches(4.8),Inches(0.4));sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
        tb(s,0.4,y+0.02,4.6,0.35,it,sz=15,b=True,c=WHITE);tb(s,0.4,y+0.45,4.7,0.5,d,sz=16,c=DARK);y+=1.1
    ib(s,5.5,1.0,4.2,y-1.1,il)
    dc={"美国":"🗽","加拿大":"🍁","墨西哥":"🌮"}.get(cn,"🌎");tb(s,8.8,4.8,1.0,0.6,dc,sz=28,c=LGRAY,a=PP_ALIGN.RIGHT)
    return s
def vs(title,bgc):
    s=ns();bg(s,bgc);tb(s,1,0.8,8,0.8,"🎬 看视频  Watch Video",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,1.8,8,0.5,title,sz=22,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
    ib(s,1.5,2.5,7,2.5,"📷 插入视频截图或粘贴视频链接");tb(s,1,5.1,8,0.3,"🔗 视频链接: ____________________",sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
    return s

n=0

# 1 COVER
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.3,8,0.8,"Global Explorer Camp",sz=36,b=True,c=TEAL,a=PP_ALIGN.CENTER)
tb(s,1,0.9,8,0.5,"环球探索沉浸式夏令营",sz=20,c=TEAL,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2),Inches(1.6),Inches(6),Inches(3.2));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=TEAL;sh.line.width=Pt(3)
tf=tb(s,2.3,1.8,5.4,2.8,"BOARDING PASS  登机牌",sz=16,b=True,c=GRAY,a=PP_ALIGN.CENTER)
ap(tf,"",sz=8);ap(tf,"Flight 航班: GR EDU-004",sz=18,c=DARK);ap(tf,"Destination 目的地:  美洲 AMERICAS",sz=20,b=True,c=TEAL)
ap(tf,"Date 日期: June 11, 2025",sz=16,c=DARK);ap(tf,"Gate 登机口: 谷雨大厅 GR EDU Hall",sz=16,c=DARK)
ap(tf,"",sz=8);ap(tf,"Passenger 旅客: 谷雨全体师生",sz=18,c=DARK);ap(tf,"",sz=6)
ap(tf,"Fasten seatbelts!  系好安全带！ ✈️",sz=14,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# 2 Schedule
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([("Session 1  上午","11:00-11:45","了解美洲 + 三个国家",TEAL),("Session 2  下午","2:00-2:45","复习总结 + 语言目标（认字写字）",BLUE),("Session 3  下午","3:00-4:30","写Booklet + 做Project",GREEN)]):
    y=0.9+i*1.5;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2));sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.1,4,0.4,nm,sz=20,b=True,c=WHITE);tb(s,0.7,y+0.5,3,0.4,tm,sz=16,c=RGBColor(0xFF,0xF3,0xE0));tb(s,4.5,y+0.15,4.8,0.9,dc,sz=18,c=WHITE)
pn(s,n)

# 3 Objectives
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.9,9,0.5,"📚 内容目标:",sz=20,b=True,c=TEAL)
tf=tb(s,0.7,1.4,8.5,1.2,"1. 了解美洲的地理位置和特点",sz=16,c=DARK);ap(tf,"2. 了解三个国家：美国、加拿大、墨西哥（国旗、首都、景点、文化）",sz=16,c=DARK)
tb(s,0.5,2.8,9,0.5,"🗣️ 语言目标:",sz=20,b=True,c=BLUE)
tb(s,0.7,3.3,4.2,0.9,"👀 我会认：美洲 美国 加拿大\n　　　　　墨西哥",sz=16,b=True,c=DARK)
tb(s,5.2,3.3,4.2,0.9,"✍️ 我会写：美洲 美国",sz=16,b=True,c=DARK)
tb(s,0.5,4.3,9,0.5,"🎨 实践目标: 完成美洲Booklet + 手工项目",sz=16,c=GREEN);pn(s,n)

# 4 S1 divider
div("Session 1  上午","了解美洲的地理位置和特点\n了解三个主要国家：美国、加拿大、墨西哥",TEAL,"🌎");n+=1

# 5 Video Americas
s=vs("认识美洲 About Americas",RGBColor(0x00,0x3D,0x35));n+=1;pn(s,n)

# 6 Where is Americas
s=ns();n+=1;bg(s,CREAM);hb(s,"🌎 美洲在哪里？ Where are the Americas?")
tb(s,0.4,0.9,4.5,0.5,"美洲分为北美洲和南美洲！",sz=20,b=True,c=TEAL);tb(s,0.4,1.4,4.5,0.4,"The Americas = North + South America!",sz=14,c=GRAY)
ib(s,5.2,0.8,4.5,4.2,"📷 世界地图\n标出美洲位置");pn(s,n)

# 7-8 About Americas
s=ns();n+=1;bg(s,CREAM);hb(s,"🌎 认识美洲  About Americas (1/2)")
for i,(t,d) in enumerate([("🌎 分为北美洲和南美洲","North America + South America"),("🌐 35个国家 35 Countries","美洲有35个国家，人口约10亿！")]):
    tb(s,0.4,0.9+i*1.2,4.8,0.5,t,sz=18,b=True,c=TEAL);tb(s,0.4,1.4+i*1.2,4.8,0.4,d,sz=15,c=DARK)
ib(s,5.5,0.9,4.2,3.5,"📷 美洲地图/风景");pn(s,n)

s=ns();n+=1;bg(s,CREAM);hb(s,"🌎 认识美洲  About Americas (2/2)")
for i,(t,d) in enumerate([("🌳 亚马逊雨林 Amazon Rainforest","世界上最大的热带雨林！"),("🏞️ 密西西比河/尼亚加拉大瀑布","Niagara Falls — 世界最壮观的瀑布之一！")]):
    tb(s,0.4,0.9+i*1.2,4.8,0.5,t,sz=18,b=True,c=TEAL);tb(s,0.4,1.4+i*1.2,4.8,0.4,d,sz=15,c=DARK)
ib(s,5.5,0.9,4.2,3.5,"📷 亚马逊雨林/尼亚加拉瀑布");pn(s,n)

# USA 美国
s=vs("🇺🇸 认识美国 About USA",RGBColor(0x1E,0x1E,0x3F));n+=1;pn(s,n)
s=cis("🇺🇸","美国","USA",USA,"国旗 + 首都",[("🏴 国旗","星条旗 Stars & Stripes 50颗星13条纹"),("🏛️ 首都","华盛顿 Washington D.C.")],"📷 美国国旗+华盛顿");n+=1;pn(s,n)
s=cis("🇺🇸","美国","USA",USA,"人口 + 语言",[("👥 人口","约3.3亿 ~330 million"),("🗣️ 语言","英语 English")],"📷 美国图片");n+=1;pn(s,n)
s=cis("🇺🇸","美国","USA",USA,"主要景点 Landmarks",[("🗽 自由女神像 Statue of Liberty","纽约标志，法国赠送的礼物！"),("🏜️ 大峡谷 Grand Canyon","世界最壮观的峡谷！"),("🎬 好莱坞 Hollywood","世界电影之都！")],"📷 自由女神/大峡谷");n+=1;pn(s,n)
s=cis("🇺🇸","美国","USA",USA,"礼节 Etiquette",[("👋 说你好","「Hello/Hi」"),("🤝 打招呼","握手或拥抱 Handshake or hug"),("🍽️ 吃饭礼节","汉堡热狗文化/小费15-20%")],"📷 美国问候/用餐");n+=1;pn(s,n)
s=cis("🇺🇸","美国","USA",USA,"美食 Food",[("🍔 汉堡 Hamburger","美国最经典的快餐！"),("🌭 热狗 Hot Dog","棒球场必备美食"),("🥧 苹果派 Apple Pie","美国传统甜点")],"📷 美国美食");n+=1;pn(s,n)

# CANADA 加拿大
s=vs("🇨🇦 认识加拿大 About Canada",RGBColor(0x6B,0x0D,0x0D));n+=1;pn(s,n)
s=cis("🇨🇦","加拿大","Canada",CANADA,"国旗 + 首都",[("🏴 国旗","红白枫叶旗 Red & White Maple Leaf 🍁"),("🏛️ 首都","渥太华 Ottawa")],"📷 加拿大国旗");n+=1;pn(s,n)
s=cis("🇨🇦","加拿大","Canada",CANADA,"人口 + 语言",[("👥 人口","约3,900万 ~39 million"),("🗣️ 语言","英语+法语 English + French")],"📷 加拿大图片");n+=1;pn(s,n)
s=cis("🇨🇦","加拿大","Canada",CANADA,"主要景点 Landmarks",[("💧 尼亚加拉瀑布 Niagara Falls","世界最著名的瀑布之一！"),("🏔️ 班夫国家公园 Banff","加拿大最美的国家公园"),("🗼 CN塔 CN Tower","多伦多标志性建筑")],"📷 尼亚加拉/班夫");n+=1;pn(s,n)
s=cis("🇨🇦","加拿大","Canada",CANADA,"礼节 Etiquette",[("👋 说你好","「Hello/Bonjour」"),("🤝 打招呼","握手 Handshake（友好但保持距离）"),("🍽️ 吃饭礼节","枫糖浆配一切/Poutine肉汁薯条")],"📷 加拿大问候");n+=1;pn(s,n)
s=cis("🇨🇦","加拿大","Canada",CANADA,"美食 Food",[("🍟 Poutine 肉汁芝士薯条","加拿大国民美食！"),("🍁 枫糖浆 Maple Syrup","加拿大特产，全球70%产量"),("🥧 Butter Tarts","经典加拿大甜点")],"📷 加拿大美食");n+=1;pn(s,n)

# MEXICO 墨西哥
s=vs("🇲🇽 认识墨西哥 About Mexico",RGBColor(0x00,0x3A,0x25));n+=1;pn(s,n)
s=cis("🇲🇽","墨西哥","Mexico",MEXICO,"国旗 + 首都",[("🏴 国旗","绿白红+鹰和蛇 Green/White/Red+Eagle & Snake"),("🏛️ 首都","墨西哥城 Mexico City")],"📷 墨西哥国旗");n+=1;pn(s,n)
s=cis("🇲🇽","墨西哥","Mexico",MEXICO,"人口 + 语言",[("👥 人口","约1.3亿 ~130 million"),("🗣️ 语言","西班牙语 Spanish")],"📷 墨西哥图片");n+=1;pn(s,n)
s=cis("🇲🇽","墨西哥","Mexico",MEXICO,"主要景点 Landmarks",[("🏛️ 奇琴伊察金字塔 Chichen Itza","玛雅文明的伟大遗迹！"),("🏖️ 坎昆海滩 Cancun","世界著名的度假胜地"),("💀 亡灵节 Day of the Dead","墨西哥最独特的文化节日")],"📷 奇琴伊察/坎昆");n+=1;pn(s,n)
s=cis("🇲🇽","墨西哥","Mexico",MEXICO,"礼节 Etiquette",[("👋 说你好","「Hola」(oh-la)"),("🤝 打招呼","拥抱+贴面 Hug + cheek kiss"),("🍽️ 吃饭礼节","玉米饼配一切/辣椒很重要")],"📷 墨西哥问候");n+=1;pn(s,n)
s=cis("🇲🇽","墨西哥","Mexico",MEXICO,"美食 Food",[("🌮 墨西哥卷饼 Taco","墨西哥最经典的美食！"),("🥑 牛油果酱 Guacamole","配玉米片的完美搭档"),("🌶️ 辣椒 Chili","墨西哥菜的灵魂！")],"📷 墨西哥美食");n+=1;pn(s,n)

# Greeting Practice
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 打招呼练习  Greeting Practice")
for i,(nm,gr,cl) in enumerate([("🇺🇸 美国","握手或拥抱+「Hello/Hi」",USA),("🇨🇦 加拿大","握手+「Hello/Bonjour」",CANADA),("🇲🇽 墨西哥","拥抱+贴面+「Hola」\n(oh-la)",MEXICO)]):
    x=0.4+i*3.2;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(2.9),Inches(3.5));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.1,2.7,0.5,nm,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER);ls=gr.split('\n');tf=tb(s,x+0.1,1.7,2.7,0.3,ls[0],sz=14,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=14,c=DARK,a=PP_ALIGN.CENTER)
    ib(s,x+0.3,2.4,2.3,1.8,"📷 动作示范")
tb(s,0.4,4.7,9,0.4,"站起来和旁边的同学练习！Stand up and practice!",sz=14,c=GRAY,a=PP_ALIGN.CENTER);pn(s,n)

# SESSION 2
div("Session 2  下午","复习总结 + 语言目标\n我会认：美洲 美国 加拿大 墨西哥\n我会写：美洲 美国",BLUE,"📖");n+=1

# Comparison blank
s=ns();n+=1;bg(s,CREAM);hb(s,"🌎 美洲三国文化对比  你知道吗？",BLUE)
tb(s,0.4,0.85,9,0.35,"你能填出来吗？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(7,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.5);t.columns[2].width=Inches(2.5);t.columns[3].width=Inches(2.6)
for r,rd in enumerate([["","🇺🇸 美国","🇨🇦 加拿大","🇲🇽 墨西哥"],["👋 说你好","？","？","？"],["🤝 打招呼","？","？","？"],["🍽️ 美食","？","？","？"],["🏛️ 首都","？","？","？"],["🏔️ 景点","？","？","？"],["🎁 代表物","？","？","？"]]):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        if r==0:rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=BLUE
        elif c==0:rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        else:rn.font.color.rgb=LGRAY;rn.font.size=Pt(20)
pn(s,n)

# Comparison answers
s=ns();n+=1;bg(s,CREAM);hb(s,"🌎 美洲三国文化对比  Comparison",BLUE)
tb(s,0.4,0.85,9,0.35,"三个国家各有特色，你最想去哪个？",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(7,4,Inches(0.3),Inches(1.25),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.8);t.columns[1].width=Inches(2.5);t.columns[2].width=Inches(2.5);t.columns[3].width=Inches(2.6)
for r,rd in enumerate([["","🇺🇸 美国","🇨🇦 加拿大","🇲🇽 墨西哥"],
    ["👋 说你好","Hello/Hi","Hello/Bonjour","Hola\n(oh-la)"],
    ["🤝 打招呼","握手或拥抱","握手(友好)","拥抱+贴面"],
    ["🍽️ 美食","汉堡Hamburger\n热狗Hot Dog","Poutine薯条\n枫糖浆","Taco卷饼\nGuacamole"],
    ["🏛️ 首都","华盛顿 D.C.","渥太华 Ottawa","墨西哥城 Mexico City"],
    ["🏔️ 景点","自由女神像\nGrand Canyon","尼亚加拉瀑布\nBanff","奇琴伊察\nCancun"],
    ["🎁 代表物","🗽 自由女神","🍁 枫叶","🌮 玉米饼"]]):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;rn=p.add_run();rn.text=ct;rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        rn.font.color.rgb=WHITE if r==0 else(DARK if c>0 else RGBColor(0x44,0x44,0x44))
        if r==0:cl.fill.solid();cl.fill.fore_color.rgb=BLUE
        elif c==0:cl.fill.solid();cl.fill.fore_color.rgb=WARM
        elif r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF5)
pn(s,n)

# Word cards 我会认
for w,py,en,sent,il in [("美洲","měi zhōu","Americas","美洲分为北美洲和南美洲。","📷 美洲地图"),("美国","měi guó","USA","美国的首都是华盛顿。","📷 自由女神像"),("加拿大","jiā ná dà","Canada","加拿大有美丽的枫叶。","📷 枫叶/班夫"),("墨西哥","mò xī gē","Mexico","墨西哥人喜欢吃玉米饼。","📷 墨西哥风景")]:
    s=ns();n+=1;bg(s,CREAM);hb(s,"👀 我会认  I Can Read",BLUE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5));sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=TEAL,a=PP_ALIGN.CENTER);tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=BLUE,a=PP_ALIGN.CENTER);ib(s,5.3,1.0,4.4,2.5,il)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2));sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=BLUE;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=BLUE);tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK);pn(s,n)

# Word games
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",BLUE)
for i,(nm,name,desc,bgc) in enumerate([("1️⃣","拍苍蝇 Fly Swatter","把字卡贴在白板上\n老师说词语，学生拍！",WARM),("2️⃣","举牌游戏 Show Me","每人4张字卡\n老师说词语，举正确的卡",RGBColor(0xE3,0xF2,0xFD)),("3️⃣","抢椅子 Musical Chairs","椅子上放字卡\n音乐停，读出词",RGBColor(0xE8,0xF5,0xE9)),("4️⃣","传话筒 Pass the Mic","传球，停下的人\n读字卡并造句",RGBColor(0xFC,0xE4,0xEC))]):
    x=0.3+i*2.4;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2));sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,nm,sz=24,a=PP_ALIGN.CENTER);tb(s,x+0.1,1.4,2.0,0.6,name,sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls=desc.split('\n');tf=tb(s,x+0.15,2.1,1.9,1.5,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=13,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.8,2.0,0.4,"低prep ✅",sz=12,b=True,c=GREEN,a=PP_ALIGN.CENTER)
pn(s,n)

# Write cards 我会写
for w,py,en,il in [("美洲","měi zhōu","Americas","📷 美洲地图"),("美国","měi guó","USA","📷 美国国旗")]:
    s=ns();n+=1;bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",BLUE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=BLUE;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=BLUE,a=PP_ALIGN.CENTER);tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,il)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8));sh2.fill.solid();sh2.fill.fore_color.rgb=RGBColor(0xE3,0xF2,0xFD);sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=BLUE);ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=BLUE);ap(tf,"1. 空中写 Air Write",sz=13,c=DARK);ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK);ap(tf,"3. 纸上写 Write 3 times",sz=13,c=DARK)
    pn(s,n)

# SESSION 3
div("Session 3  下午","写Booklet + 做Project\n3:00 - 4:30",GREEN,"🎨");n+=1

# Booklet
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成"探索美洲"练习册',GREEN);ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图");pn(s,n)

# Project overview
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 Project Time!  4个手工项目",GREEN)
for i,(pr,nm,md,bc) in enumerate([("PROJECT 1","🧩 美洲拼图","group project",WARM),("PROJECT 2","🏛️ 地标手工","独立完成",RGBColor(0xE3,0xF2,0xFD)),("PROJECT 3","🍫 巧克力装饰","独立完成",RGBColor(0xE8,0xF5,0xE9)),("PROJECT 4","💀 亡灵节面具","group project",RGBColor(0xFC,0xE4,0xEC))]):
    x=0.3+i*2.4;sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2));sh.fill.solid();sh.fill.fore_color.rgb=bc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.3,pr,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER);tb(s,x+0.1,1.3,2.0,0.6,nm,sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.9,2.0,0.3,f"({md})",sz=11,c=GRAY,a=PP_ALIGN.CENTER);ib(s,x+0.2,2.4,1.8,2.3,"📷 示范")
pn(s,n)

# Individual projects
for pr,nm,md,bc,cl in [("PROJECT 1","🧩 美洲拼图  Americas Puzzle Map","group project",WARM,GREEN),("PROJECT 2","🏛️ 地标手工  Landmark Craft","独立完成",RGBColor(0xE3,0xF2,0xFD),BLUE),("PROJECT 3","🍫 巧克力装饰  Chocolate Art","独立完成",RGBColor(0xE8,0xF5,0xE9),GREEN),("PROJECT 4","💀 亡灵节面具  Day of Dead Mask","group project",RGBColor(0xFC,0xE4,0xEC),RGBColor(0xC2,0x18,0x5B))]:
    s=ns();n+=1;bg(s,bc);hb(s,f"{pr}: {nm}",cl);tb(s,0.4,0.85,9,0.3,f"({md})",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,0.4,1.3,4.4,3.5,"📷 示范图片/截图");ib(s,5.2,1.3,4.5,3.5,"🎬 教学视频/步骤视频");pn(s,n)

# Visa stamp
s=ns();n+=1;bg(s,CREAM);tb(s,1,0.5,8,0.8,"🪪 美洲签证章  Americas Visa Stamp",sz=30,b=True,c=TEAL,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.5),Inches(3),Inches(3));sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=TEAL;sh.line.width=Pt(5)
tf=tb(s,3.6,1.8,2.8,2.5,"AMERICAS\n美洲",sz=28,b=True,c=TEAL,a=PP_ALIGN.CENTER);ap(tf,"✓ VISITED",sz=16,b=True,c=GREEN,a=PP_ALIGN.CENTER);ap(tf,"6/11/2025",sz=12,c=GRAY,a=PP_ALIGN.CENTER);ap(tf,"美国 · 加拿大 · 墨西哥",sz=12,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1,4.7,8,0.4,"恭喜你完成美洲之旅！Congratulations! 🎉",sz=16,b=True,c=TEAL,a=PP_ALIGN.CENTER);pn(s,n)

# Tomorrow
s=ns();n+=1;bg(s,TEAL)
tb(s,1,1.0,8,0.8,"✈️ 明天航班  Tomorrow's Flight",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,2,2.2,6,2.5,"Day 5 展览日 Exhibition Day!",sz=22,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"🎉 明天是展览日！",sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER);ap(tf,"",sz=10)
ap(tf,"展示你的护照和所有手工项目！",sz=20,c=WHITE,a=PP_ALIGN.CENTER);ap(tf,"Show your passport and all your projects!",sz=18,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER)
ap(tf,"",sz=10);ap(tf,"See you tomorrow, explorers! 明天见！",sz=16,c=RGBColor(0xFF,0xF3,0xE0),a=PP_ALIGN.CENTER);pn(s,n)

OUT='/Users/Huan/projects/summercourse/Chinese/世界旅行world_trip_pbl/day4_americas_v2.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
