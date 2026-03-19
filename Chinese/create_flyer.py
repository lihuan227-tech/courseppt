#!/usr/bin/env python3
"""Generate a one-page Chinese course flyer in PPT format."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- constants ---
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
DARK_ORANGE = RGBColor(0xE6, 0x7E, 0x00)
GREEN = RGBColor(0x9D, 0xC6, 0x4D)
DARK_GREEN = RGBColor(0x5E, 0x94, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xFF, 0xF8, 0xF0)
LIGHT_GREEN_BG = RGBColor(0xF1, 0xF8, 0xF3)

prs = Presentation()
# Letter size (8.5 x 11 inches)
prs.slide_width = Inches(8.5)
prs.slide_height = Inches(11)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)


def add_shape(left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=10, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='KaiTi'):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = font_name
    return p


def add_para(tf, text, font_size=10, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(0), font_name='KaiTi'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    p.font.name = font_name
    return p


# ============================
# HEADER BANNER (green background)
# ============================
banner = add_shape(Inches(0), Inches(0), Inches(8.5), Inches(1.4), fill_color=GREEN)
# Orange border at bottom
border_line = add_shape(Inches(0), Inches(1.36), Inches(8.5), Inches(0.04), fill_color=ORANGE)

txBox = add_textbox(Inches(0.5), Inches(0.15), Inches(7.5), Inches(0.6))
tf = txBox.text_frame
set_text(tf, '暑假中文课程', font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox2 = add_textbox(Inches(0.5), Inches(0.7), Inches(7.5), Inches(0.35))
tf2 = txBox2.text_frame
set_text(tf2, 'Summer Chinese Course', font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox3 = add_textbox(Inches(0.5), Inches(1.05), Inches(7.5), Inches(0.28))
tf3 = txBox3.text_frame
set_text(tf3, '谷雨中文 · 2025 Summer Program', font_size=14, color=RGBColor(0xFF, 0xFF, 0xE0), alignment=PP_ALIGN.CENTER)

# ============================
# FEATURES SECTION
# ============================
y_feat = Inches(1.55)

txBox = add_textbox(Inches(0.5), y_feat, Inches(7.5), Inches(0.35))
tf = txBox.text_frame
set_text(tf, '课程特色', font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

features = [
    ('✏️ 零基础入门', '从握笔、笔画、基础汉字学起，认一认、读一读、写一写，轻松开启中文学习。'),
    ('🔍 查漏补缺', '针对学过中文但基础不牢的孩子，利用暑假复习巩固，补齐短板，迎接新学期更轻松。'),
    ('🚀 进阶提升', '加强识字、阅读、写字与表达训练，全面提升中文能力，让孩子更上一层楼。'),
]

card_w = Inches(2.3)
card_h = Inches(1.3)
gap = Inches(0.2)
start_x = Inches(0.75)
y_cards = y_feat + Inches(0.4)

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    card = add_shape(x, y_cards, card_w, card_h, fill_color=WHITE, border_color=ORANGE, border_width=Pt(2))

    txBox = add_textbox(x + Inches(0.1), y_cards + Inches(0.08), card_w - Inches(0.2), Inches(0.3))
    tf = txBox.text_frame
    set_text(tf, title, font_size=14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    txBox = add_textbox(x + Inches(0.1), y_cards + Inches(0.4), card_w - Inches(0.2), Inches(0.85))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, desc, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# LEVELS SECTION
# ============================
y_levels = y_cards + card_h + Inches(0.25)

txBox = add_textbox(Inches(0.5), y_levels, Inches(7.5), Inches(0.35))
tf = txBox.text_frame
set_text(tf, '课程介绍', font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

levels = [
    {
        'name': 'Level 1',
        'target': '零基础，希望在暑假进行中文入门，建立基本识字、写字和阅读兴趣的学生',
        'content': '基础识字、写字、分级阅读，培养阅读兴趣和信心',
        'assessment': '无需测评',
    },
    {
        'name': 'Level 2',
        'target': '已有一定中文学习基础（半年到一年），希望转入部编版体系或巩固提升',
        'content': '《部编版语文一年级上册》重难点课文',
        'assessment': '测评正确率 < 70% 建议选择本级别',
    },
    {
        'name': 'Level 3',
        'target': '已有中文基础（一年到一年半），希望转入部编版体系或巩固提升',
        'content': '《部编版语文一年级下册》重难点课文',
        'assessment': '测评正确率 < 70% 建议选择本级别',
    },
    {
        'name': 'Level 4',
        'target': '已有中文基础（一年半到两年左右），希望转入部编版体系或巩固提升',
        'content': '《部编版语文二年级上册》重难点课文',
        'assessment': '测评正确率 < 70% 建议选择本级别',
    },
    {
        'name': 'Level 5',
        'target': '已有中文基础（两到三年左右），希望转入部编版体系或巩固提升',
        'content': '《部编版语文二年级下册》重难点课文',
        'assessment': '测评正确率 < 70% 建议选择本级别',
    },
]

# Table-like layout for levels
y_table = y_levels + Inches(0.4)
table_w = Inches(7.2)
table_x = Inches(0.65)

# Header row
header_bg = add_shape(table_x, y_table, table_w, Inches(0.35), fill_color=ORANGE)
cols = [Inches(0.8), Inches(2.8), Inches(2.2), Inches(1.4)]
col_starts = [table_x]
for c in cols[:-1]:
    col_starts.append(col_starts[-1] + c)
headers = ['级别', '🎯 适合学生', '📖 教学内容', '📋 测评']
for i, (header, cw) in enumerate(zip(headers, cols)):
    txBox = add_textbox(col_starts[i] + Inches(0.05), y_table + Inches(0.04), cw - Inches(0.1), Inches(0.28))
    tf = txBox.text_frame
    set_text(tf, header, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Data rows
for j, level in enumerate(levels):
    y_row = y_table + Inches(0.35) + j * Inches(0.72)
    bg_color = LIGHT_BG if j % 2 == 0 else WHITE
    row_bg = add_shape(table_x, y_row, table_w, Inches(0.7), fill_color=bg_color, border_color=RGBColor(0xE0, 0xE0, 0xE0), border_width=Pt(0.5))

    data = [level['name'], level['target'], level['content'], level['assessment']]
    for i, (text, cw) in enumerate(zip(data, cols)):
        txBox = add_textbox(col_starts[i] + Inches(0.05), y_row + Inches(0.05), cw - Inches(0.1), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        if i == 0:
            set_text(tf, text, font_size=12, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
        else:
            set_text(tf, text, font_size=10, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)

# ============================
# SCHEDULE / SESSION INFO
# ============================
y_sched = y_table + Inches(0.35) + 5 * Inches(0.72) + Inches(0.2)

txBox = add_textbox(Inches(0.5), y_sched, Inches(7.5), Inches(0.35))
tf = txBox.text_frame
set_text(tf, '课程安排', font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

y_info = y_sched + Inches(0.38)
info_bg = add_shape(Inches(0.65), y_info, Inches(7.2), Inches(0.32), fill_color=LIGHT_GREEN_BG)
txBox = add_textbox(Inches(0.75), y_info + Inches(0.04), Inches(7.0), Inches(0.26))
tf = txBox.text_frame
set_text(tf, '📅 每个 Session 共 4 节课（每周 2 次 × 2 周）  |  5 个级别  |  每级别 4 个 Session 可选', font_size=11, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)

# Session schedule summary
y_sess = y_info + Inches(0.42)
sess_data = [
    ('Session 1', '6/8–6/19'),
    ('Session 2', '6/22–7/10'),
    ('Session 3', '7/13–7/24'),
    ('Session 4', '7/27–8/7'),
]

sess_card_w = Inches(1.65)
sess_gap = Inches(0.15)
sess_start_x = Inches(0.75)

for i, (sess_name, dates) in enumerate(sess_data):
    x = sess_start_x + i * (sess_card_w + sess_gap)
    card = add_shape(x, y_sess, sess_card_w, Inches(0.6), fill_color=WHITE, border_color=DARK_GREEN, border_width=Pt(2))

    txBox = add_textbox(x + Inches(0.05), y_sess + Inches(0.05), sess_card_w - Inches(0.1), Inches(0.25))
    tf = txBox.text_frame
    set_text(tf, sess_name, font_size=14, bold=True, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)

    txBox = add_textbox(x + Inches(0.05), y_sess + Inches(0.32), sess_card_w - Inches(0.1), Inches(0.22))
    tf = txBox.text_frame
    set_text(tf, dates, font_size=11, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# CLASS FORMAT
# ============================
y_format = y_sess + Inches(0.75)

format_bg = add_shape(Inches(0.65), y_format, Inches(7.2), Inches(1.0), fill_color=WHITE, border_color=ORANGE, border_width=Pt(2))

txBox = add_textbox(Inches(0.85), y_format + Inches(0.06), Inches(6.8), Inches(0.3))
tf = txBox.text_frame
set_text(tf, '📋 课程详情', font_size=14, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

details = [
    '上课形式：线上小班（Zoom）  |  每班 4–8 人',
    '课时安排：每节课 50 分钟  |  每周 2 次  |  每 Session 共 4 节课',
    '适合年龄：K–5 年级（按中文水平分级，非按年级分班）',
]

for k, detail in enumerate(details):
    txBox = add_textbox(Inches(1.0), y_format + Inches(0.36) + k * Inches(0.2), Inches(6.5), Inches(0.22))
    tf = txBox.text_frame
    set_text(tf, detail, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# FOOTER / CTA
# ============================
y_footer = y_format + Inches(1.15)

footer_bg = add_shape(Inches(0), y_footer, Inches(8.5), Inches(0.6), fill_color=ORANGE)

txBox = add_textbox(Inches(0.5), y_footer + Inches(0.06), Inches(7.5), Inches(0.3))
tf = txBox.text_frame
set_text(tf, '立即报名  Register Now', font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox = add_textbox(Inches(0.5), y_footer + Inches(0.35), Inches(7.5), Inches(0.22))
tf = txBox.text_frame
set_text(tf, '谷雨中文  |  www.mygredu.com', font_size=12, color=RGBColor(0xFF, 0xFF, 0xE0), alignment=PP_ALIGN.CENTER)

# Save
output_path = '/Users/Huan/projects/summercourse/Chinese/chinese_course_flyer.pptx'
prs.save(output_path)
print(f'Flyer saved to: {output_path}')
