#!/usr/bin/env python3
"""Generate a one-page Chinese course flyer in PPT format (v2 - compact with all info)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- constants ---
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
DARK_ORANGE = RGBColor(0xE6, 0x7E, 0x00)
GREEN = RGBColor(0x9D, 0xC6, 0x4D)
DARK_GREEN = RGBColor(0x5E, 0x94, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xFF, 0xF8, 0xF0)
LIGHT_GREEN_BG = RGBColor(0xF1, 0xF8, 0xF3)
BLUE_LINK = RGBColor(0x19, 0x76, 0xD2)

prs = Presentation()
prs.slide_width = Inches(8.5)
prs.slide_height = Inches(11)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_shape(left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=10, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='KaiTi'):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = font_name
    return p


def add_para(tf, text, font_size=10, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(0), font_name='KaiTi'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    p.font.name = font_name
    return p


# ============================
# 1. HEADER BANNER
# ============================
add_rect(Inches(0), Inches(0), Inches(8.5), Inches(1.05), fill_color=GREEN)
add_rect(Inches(0), Inches(1.02), Inches(8.5), Inches(0.03), fill_color=ORANGE)

txBox = add_textbox(Inches(0.5), Inches(0.1), Inches(7.5), Inches(0.5))
set_text(txBox.text_frame, '暑假中文课程 Summer Chinese Course', font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox = add_textbox(Inches(0.5), Inches(0.6), Inches(7.5), Inches(0.35))
set_text(txBox.text_frame, '谷雨中文 GR EDU · 2025 Summer · Onsite Classes · K–5年级', font_size=13, color=RGBColor(0xFF, 0xFF, 0xE0), alignment=PP_ALIGN.CENTER)

# ============================
# 2. 课程特色 FEATURES (compact 1-row)
# ============================
y = Inches(1.18)

txBox = add_textbox(Inches(0.4), y, Inches(7.7), Inches(0.3))
set_text(txBox.text_frame, '课程特色', font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

features = [
    ('✏️ 零基础入门', '从笔画、汉字学起\n轻松开启中文学习'),
    ('🔍 查漏补缺', '暑假复习巩固\n补齐短板迎新学期'),
    ('🚀 进阶提升', '加强识字阅读写字\n全面提升中文能力'),
]

card_w = Inches(2.3)
card_h = Inches(0.85)
gap = Inches(0.15)
start_x = Inches(0.7)
y_cards = y + Inches(0.32)

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    add_shape(x, y_cards, card_w, card_h, fill_color=WHITE, border_color=ORANGE, border_width=Pt(2))

    txBox = add_textbox(x + Inches(0.08), y_cards + Inches(0.05), card_w - Inches(0.16), Inches(0.22))
    set_text(txBox.text_frame, title, font_size=13, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    txBox = add_textbox(x + Inches(0.08), y_cards + Inches(0.28), card_w - Inches(0.16), Inches(0.55))
    txBox.text_frame.word_wrap = True
    set_text(txBox.text_frame, desc, font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# 3. 课程级别 LEVELS TABLE (with QR code space)
# ============================
y_levels = y_cards + card_h + Inches(0.18)

txBox = add_textbox(Inches(0.4), y_levels, Inches(5.5), Inches(0.3))
set_text(txBox.text_frame, '课程级别与上课时间', font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)

# QR code placeholder label
txBox = add_textbox(Inches(6.0), y_levels - Inches(0.02), Inches(2.2), Inches(0.3))
set_text(txBox.text_frame, '扫码测评选级 →', font_size=11, bold=True, color=BLUE_LINK, alignment=PP_ALIGN.CENTER)

levels = [
    {'name': 'L1 零基础', 'desc': '基础识字、写字、分级阅读', 'sched': 'Tue/Thu 5:00–6:00', 'grade': 'K–1st', 'test': '无需测评'},
    {'name': 'L2', 'desc': '部编版语文一年级上册', 'sched': 'Mon/Wed 4:00–5:00', 'grade': '1st–2nd', 'test': '测评选级'},
    {'name': 'L3', 'desc': '部编版语文一年级下册', 'sched': 'Mon/Wed 5:00–6:00', 'grade': '2nd–3rd', 'test': '测评选级'},
    {'name': 'L4', 'desc': '部编版语文二年级上册', 'sched': 'Fri 4:00–6:00', 'grade': '3rd–5th', 'test': '测评选级'},
    {'name': 'L5', 'desc': '部编版语文二年级下册', 'sched': 'Tue/Thu 4:00–5:00', 'grade': '4th+', 'test': '测评选级'},
]

y_table = y_levels + Inches(0.32)
table_x = Inches(0.4)
table_w = Inches(5.5)
row_h = Inches(0.48)

# Header
add_rect(table_x, y_table, table_w, Inches(0.32), fill_color=ORANGE)
cols = [Inches(1.0), Inches(2.0), Inches(1.5), Inches(1.0)]
col_starts = [table_x]
for c in cols[:-1]:
    col_starts.append(col_starts[-1] + c)
headers = ['级别', '教学内容', '上课时间', '参考年级']
for i, (h, cw) in enumerate(zip(headers, cols)):
    txBox = add_textbox(col_starts[i] + Inches(0.03), y_table + Inches(0.04), cw - Inches(0.06), Inches(0.25))
    set_text(txBox.text_frame, h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Data rows
for j, lv in enumerate(levels):
    y_row = y_table + Inches(0.32) + j * row_h
    bg = LIGHT_BG if j % 2 == 0 else WHITE
    add_rect(table_x, y_row, table_w, row_h, fill_color=bg)

    data = [lv['name'], lv['desc'], lv['sched'], lv['grade']]
    for i, (text, cw) in enumerate(zip(data, cols)):
        txBox = add_textbox(col_starts[i] + Inches(0.03), y_row + Inches(0.06), cw - Inches(0.06), row_h - Inches(0.1))
        txBox.text_frame.word_wrap = True
        if i == 0:
            set_text(txBox.text_frame, text, font_size=10, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
        else:
            set_text(txBox.text_frame, text, font_size=9, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# --- QR Code placeholder boxes (right side, next to table) ---
qr_x = Inches(6.1)
qr_size = Inches(1.1)
qr_gap = Inches(0.12)

# Main QR placeholder (for general assessment)
qr_y = y_table + Inches(0.05)
add_shape(qr_x, qr_y, qr_size, qr_size, border_color=ORANGE, border_width=Pt(2))
txBox = add_textbox(qr_x + Inches(0.05), qr_y + Inches(0.3), qr_size - Inches(0.1), Inches(0.5))
set_text(txBox.text_frame, '测评\nQR Code', font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Second QR placeholder (for registration)
qr_y2 = qr_y + qr_size + qr_gap
add_shape(qr_x, qr_y2, qr_size, qr_size, border_color=DARK_GREEN, border_width=Pt(2))
txBox = add_textbox(qr_x + Inches(0.05), qr_y2 + Inches(0.3), qr_size - Inches(0.1), Inches(0.5))
set_text(txBox.text_frame, '报名\nQR Code', font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Label under QR codes
txBox = add_textbox(Inches(5.95), qr_y2 + qr_size + Inches(0.04), Inches(1.4), Inches(0.2))
set_text(txBox.text_frame, '* 按中文水平分级，非按年级', font_size=7, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# 4. 上课时间 SESSION DATES
# ============================
y_sess_section = y_table + Inches(0.32) + 5 * row_h + Inches(0.18)

txBox = add_textbox(Inches(0.4), y_sess_section, Inches(7.7), Inches(0.3))
set_text(txBox.text_frame, 'Session 时间安排', font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

y_sess = y_sess_section + Inches(0.32)
add_shape(Inches(0.4), y_sess, Inches(7.7), Inches(0.28), fill_color=LIGHT_GREEN_BG)
txBox = add_textbox(Inches(0.5), y_sess + Inches(0.03), Inches(7.5), Inches(0.22))
set_text(txBox.text_frame, '每 Session 共 4 节课（每周 2 次 × 2 周）· 5 个级别 · 可按 Session 单独报名', font_size=10, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)

sess_data = [
    ('Session 1', 'Wk 1: 6/8–6/12\nWk 2: 6/15–6/19'),
    ('Session 2', 'Wk 3: 6/22–6/26\nWk 4: 7/6–7/10'),
    ('Session 3', 'Wk 5: 7/13–7/17\nWk 6: 7/20–7/24'),
    ('Session 4', 'Wk 7: 7/27–7/31\nWk 8: 8/3–8/7'),
]

y_sess_cards = y_sess + Inches(0.36)
sess_w = Inches(1.82)
sess_gap = Inches(0.12)
sess_start = Inches(0.48)

for i, (name, dates) in enumerate(sess_data):
    x = sess_start + i * (sess_w + sess_gap)
    add_shape(x, y_sess_cards, sess_w, Inches(0.72), fill_color=WHITE, border_color=DARK_GREEN, border_width=Pt(1.5))

    txBox = add_textbox(x + Inches(0.05), y_sess_cards + Inches(0.04), sess_w - Inches(0.1), Inches(0.2))
    set_text(txBox.text_frame, name, font_size=12, bold=True, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)

    txBox = add_textbox(x + Inches(0.05), y_sess_cards + Inches(0.26), sess_w - Inches(0.1), Inches(0.44))
    txBox.text_frame.word_wrap = True
    set_text(txBox.text_frame, dates, font_size=8.5, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# 5. 学费 TUITION
# ============================
y_tuition = y_sess_cards + Inches(0.85)

txBox = add_textbox(Inches(0.4), y_tuition, Inches(7.7), Inches(0.3))
set_text(txBox.text_frame, '学费 Tuition', font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

y_tuit_box = y_tuition + Inches(0.3)
add_shape(Inches(0.9), y_tuit_box, Inches(6.7), Inches(0.7), fill_color=WHITE, border_color=ORANGE, border_width=Pt(2))

# Tuition content
txBox = add_textbox(Inches(1.1), y_tuit_box + Inches(0.08), Inches(3.0), Inches(0.55))
tf = txBox.text_frame
set_text(tf, '中文课 · $60/week · $120/session', font_size=12, bold=True, color=DARK_ORANGE, alignment=PP_ALIGN.LEFT)
add_para(tf, 'Onsite 小班授课 · 6–12人 · 每周两节课', font_size=9, color=GRAY_TEXT, space_before=Pt(4))

txBox = add_textbox(Inches(4.5), y_tuit_box + Inches(0.08), Inches(2.8), Inches(0.55))
tf = txBox.text_frame
set_text(tf, '缴费方式 Payment: Zelle', font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)
add_para(tf, 'gredu2019@gmail.com', font_size=10, bold=True, color=BLUE_LINK, space_before=Pt(2))
add_para(tf, '备注：学生姓名 + Elective', font_size=8, color=GRAY_TEXT, space_before=Pt(2))

# ============================
# 6. FOOTER / CTA
# ============================
y_footer = y_tuit_box + Inches(0.82)

add_rect(Inches(0), y_footer, Inches(8.5), Inches(0.55), fill_color=ORANGE)

txBox = add_textbox(Inches(0.5), y_footer + Inches(0.05), Inches(7.5), Inches(0.25))
set_text(txBox.text_frame, '立即报名 Register Now', font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox = add_textbox(Inches(0.5), y_footer + Inches(0.3), Inches(7.5), Inches(0.2))
set_text(txBox.text_frame, '谷雨中文 GR EDU  |  www.mygredu.com/summer-camp/summer-electives', font_size=10, color=RGBColor(0xFF, 0xFF, 0xE0), alignment=PP_ALIGN.CENTER)

# Save
output_path = '/Users/Huan/projects/summercourse/Chinese/chinese_flyer_v2.pptx'
prs.save(output_path)
print(f'Flyer saved to: {output_path}')
