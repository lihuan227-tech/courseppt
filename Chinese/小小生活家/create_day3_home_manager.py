#!/usr/bin/env python3
"""
小小生活家 Little Homemaker — Day 3: 小小家庭管家 Little Home Helper
Built to the detailed lesson plan:
  Part 1 爸爸妈妈忙碌的一天 (hook: observe + brainstorm + 我能帮什么忙)
  Part 2 认识各种家务 (什么是家务 + 分类 + 动作猜家务)
  Part 3 技能一 洗碗
  Part 4 技能二 洗衣服 (洗衣机 + 检查口袋 + 旅行手洗)
  Part 5 技能三 叠衣服和收纳 (竖式叠衣法 + 抽屉整理)
  Part 6 技能四 擦桌子和扫地 + 任务顺序挑战
  下午   语言目标 我会认(5)/我会写(4)/我会说 + 四个技能轮换站 + 家务打卡表
Palette: Cozy Home (plum + amber) — distinct from Day 1 green/orange and Day 2 teal/coral.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Cozy Home (温馨小家) ---
PLUM = RGBColor(0x7B,0x4F,0x9E)   # primary: friendly plum purple
AMBER = RGBColor(0xE8,0x92,0x2E)  # accent: warm amber
DEEP = RGBColor(0x45,0x32,0x5E)   # deep plum (dark slides / dividers)
SUNNY = RGBColor(0xF4,0xC1,0x3D)
CREAM = RGBColor(0xFA,0xF5,0xEE)  # background cream
WARM = RGBColor(0xF6,0xEE,0xF9)   # soft lilac panel
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
IMGBG = RGBColor(0xEC,0xEC,0xE6)
GREEN_OK = RGBColor(0x2E,0x9E,0x7A)
RED = RGBColor(0xD8,0x45,0x3A)
GOLD = RGBColor(0xD1,0x8F,0x0A)
BLUE = RGBColor(0x3E,0x8E,0xC4)
TEAL = RGBColor(0x16,0x84,0x8A)
CORAL = RGBColor(0xEF,0x6B,0x53)
PINKBG = RGBColor(0xFD,0xEF,0xE6)

# 四色技能站 station colors
STA1 = PLUM     # 衣服收纳站
STA2 = BLUE     # 书桌整理站
STA3 = CORAL    # 厨房清洁站
STA4 = TEAL     # 洗衣练习站

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True
    lines=txt.split("\n")
    p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=lines[0];r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
    for ln in lines[1:]:ap(tf,ln,sz=sz,b=b,c=c,a=a)
    return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    for ln in txt.split("\n"):
        p=tf.add_paragraph()
        if a:p.alignment=a
        r=p.add_run();r.text=ln;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=PLUM,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n):
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(9.18),Inches(5.28),Inches(0.5),Inches(0.30))
    chip.fill.solid();chip.fill.fore_color.rgb=WHITE;chip.line.color.rgb=LGRAY;chip.line.width=Pt(0.75)
    bx=s.shapes.add_textbox(Inches(9.18),Inches(5.30),Inches(0.5),Inches(0.26));tf=bx.text_frame;tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    r=p.add_run();r.text=str(n);r.font.size=Pt(10);r.font.color.rgb=GRAY;r.font.name='KaiTi'
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    lines=sub.split("\n")
    tf=tb(s,0.4,2.75,9.2,1.6,lines[0],sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    for ln in lines[1:]:ap(tf,ln,sz=20,c=WHITE,a=PP_ALIGN.CENTER)
    return s
def step_head(s,num,cn,en,color):
    bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.15),Inches(9.4),Inches(0.72))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.45),Inches(0.24),Inches(0.54),Inches(0.54))
    circ.fill.solid();circ.fill.fore_color.rgb=WHITE;circ.line.fill.background()
    tb(s,0.45,0.28,0.54,0.46,num,sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,1.2,0.19,8.2,0.42,cn,sz=21,b=True,c=WHITE)
    tb(s,1.22,0.60,8.2,0.24,en,sz=10,c=WHITE)
def warn(s,y,txt,w=9.4,x=0.3):
    """统一红色安全警示条"""
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(0.55))
    sh.fill.solid();sh.fill.fore_color.rgb=RED;sh.line.fill.background()
    tb(s,x+0.2,y+0.08,w-0.4,0.4,f"⚠️ 安全提醒：{txt}",sz=13,b=True,c=WHITE)

n=0

# ============================================================
# 1 COVER
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.25,8,0.7,"Little Home Helper",sz=32,b=True,c=PLUM,a=PP_ALIGN.CENTER)
tb(s,1,0.85,8,0.45,"小小家庭管家",sz=20,c=PLUM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.5),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=PLUM;sh.line.color.rgb=AMBER;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.8),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=AMBER;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.05,2.8,0.4,"DAY 3",sz=16,b=True,c=AMBER,a=PP_ALIGN.CENTER)
ap(tf,"🏠",sz=48,a=PP_ALIGN.CENTER)
ap(tf,"小小家庭管家",sz=18,b=True,c=PLUM,a=PP_ALIGN.CENTER)
ap(tf,"LITTLE HOME HELPER",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,5.05,8,0.4,"💪 我会做家务，我能帮助家人！I can help my family!",sz=14,b=True,c=AMBER,a=PP_ALIGN.CENTER)
notes(s,"封面：欢迎学生。今天的任务——学习成为小小家庭管家，完成四个家务技能站打卡，获得家庭管家徽章，并把学会的技能带回家实践。")
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日安排  Today's Plan")
for i,(nm,tm,dc,cl) in enumerate([
    ("上午 · 认识家务","Part 1–2","爸爸妈妈忙碌的一天 → 什么是家务 → 家务分类",PLUM),
    ("上午 · 四大技能","Part 3–6","洗碗 · 洗衣服 · 叠衣收纳 · 擦桌子扫地",AMBER),
    ("下午 · 语言 + 挑战","技能站","我会认/写/说 + 四站轮换挑战 + 家务打卡表",DEEP)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,5,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.62,3,0.4,tm,sz=14,c=WARM)
    tb(s,4.5,y+0.38,5.1,0.6,dc,sz=13,c=WHITE)
notes(s,"上午约2.5–3小时：认识家务+四大技能。下午约2–2.5小时：语言学习+四个技能轮换站+家务打卡表。")
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.85,9,0.5,"🏠 内容目标  Content:",sz=19,b=True,c=PLUM)
tf=tb(s,0.7,1.32,9,1.5,"1. 认识家庭中常见的家务，了解爸爸妈妈每天有多辛苦",sz=14,c=DARK)
ap(tf,"2. 学会四大技能：洗碗 · 洗衣服(机洗+手洗) · 叠衣服和收纳 · 擦桌子扫地",sz=14,c=DARK)
ap(tf,"3. 学会安排任务顺序：先做什么，再做什么，最后做什么",sz=14,c=DARK)
ap(tf,"4. 完成四个技能站打卡，制作「家务打卡表」带回家实践",sz=14,c=DARK)
tb(s,0.5,3.0,9,0.5,"🗣️ 语言目标  Language:",sz=19,b=True,c=AMBER)
tb(s,0.7,3.5,5.2,0.5,"👀 我会认：洗 叠 碗 衣服 扫地",sz=14,b=True,c=DARK)
tb(s,5.9,3.5,3.8,0.5,"✍️ 我会写：洗 衣 扫 地",sz=14,b=True,c=DARK)
tb(s,0.5,4.25,9,0.5,"🎖️ 实践目标：完成四个家务技能站打卡，成为小小家庭管家，获得勋章！",sz=14,c=DEEP)
notes(s,"介绍课程主线：爸爸妈妈每天很忙，小朋友要学习成为小小家庭管家，完成家务技能训练，获得家庭管家徽章，并把学会的技能带回家实践。")
pn(s,n)

# ============================================================
# 4 PART 1 DIVIDER
# ============================================================
div("Part 1 · 上午","爸爸妈妈忙碌的一天\n👀 观察图片  💭 头脑风暴  🙋 我能帮什么忙?",PLUM,"👨‍👩‍👧")
n+=1

# ============================================================
# 5 观察图片 — 家长忙碌的一天
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👀 看一看：爸爸妈妈的一天  A Busy Day",AMBER)
tb(s,0.4,0.85,9.2,0.32,"仔细观察这张图 — 先别说答案，看看你能发现什么！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.3,1.25,4.9,3.55,"📷 家长忙碌的一天\n上班 · 做饭 · 买菜 · 接送\n洗衣 · 洗碗 · 打扫 · 检查作业")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.4),Inches(1.25),Inches(4.3),Inches(3.55))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=PLUM;panel.line.width=Pt(2.5)
hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.4),Inches(1.25),Inches(4.3),Inches(0.5))
hd.fill.solid();hd.fill.fore_color.rgb=PLUM;hd.line.fill.background()
tb(s,5.55,1.32,4.0,0.4,"❓ 说一说  Discuss",sz=14,b=True,c=WHITE)
tf=tb(s,5.55,1.9,4.0,2.7,"· 你在图片里看到了什么?",sz=14,c=DARK)
for q in ["· 爸爸妈妈每天要做哪些事情?","· 哪些事情发生在家里?","· 哪些事情发生在外面?"]:
    ap(tf,"",sz=13);ap(tf,q,sz=14,c=DARK)
tb(s,5.55,4.35,4.0,0.35,"🙋 先和同桌小声说一说！",sz=12,b=True,c=AMBER)
notes(s,"情境导入(5分钟)：展示一位家长忙碌的一天(上班/做饭/买菜/接送孩子/洗衣服/洗碗/打扫卫生/检查作业)。不要直接把答案写出来，先让学生观察和猜。提问4个问题，引导区分家里/外面的事情。")
pn(s,n)

# ============================================================
# 6 Brainstorm — 爸爸妈妈每天还要做什么?
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💭 头脑风暴：爸爸妈妈每天还要做什么?",PLUM)
tb(s,0.4,0.85,9.2,0.32,"学生先说，老师再点击 — 答案一个一个出现！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
answers=[("💼","上班"),("🍳","做饭"),("🛒","买菜"),("👕","洗衣服"),("🍽️","洗碗"),
         ("🧹","打扫卫生"),("🚗","接送孩子"),("📖","检查作业"),("🛏️","整理房间"),("❤️","照顾家人")]
for i,(em,txt) in enumerate(answers):
    col=i%5;row=i//5
    x=0.4+col*1.9;y=1.35+row*1.35
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(1.7),Inches(1.15))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=AMBER;card.line.width=Pt(2)
    tb(s,x,y+0.1,1.7,0.55,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.68,1.7,0.4,txt,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
disc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.25),Inches(9.4),Inches(0.85))
disc.fill.solid();disc.fill.fore_color.rgb=WARM;disc.line.color.rgb=PLUM;disc.line.width=Pt(2)
tb(s,0.5,4.33,9.0,0.35,"🤔 一起讨论  Think together:",sz=13,b=True,c=PLUM)
tb(s,0.5,4.68,9.0,0.35,"爸爸妈妈一天要做这么多事情，他们会不会累？我们可以怎么帮忙？",sz=14,b=True,c=DARK)
notes(s,"头脑风暴(5分钟)：先让学生自由回答，老师每听到一个就点击出现对应卡片(放映时建议给每张卡片加「点击出现」动画)。答案：上班/做饭/买菜/洗衣服/洗碗/打扫卫生/接送孩子/检查作业/整理房间/照顾家人。最后讨论：爸爸妈妈会不会累？引出「我们可以帮忙」。")
pn(s,n)

# ============================================================
# 7 我能帮什么忙? — 红黄绿判断
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🙋 我能帮什么忙?  What Can I Do?",AMBER)
tb(s,0.4,0.85,9.2,0.3,"看到一件家务，想一想：它属于哪种颜色?",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
cols=[("🟢 我可以自己做",GREEN_OK,["叠衣服","收玩具","擦桌子","扫地","洗自己的袜子"]),
      ("🟡 需要大人帮忙",GOLD,["洗碗","使用洗衣机"]),
      ("🔴 对小朋友不安全",RED,["使用炉灶","使用锋利的刀","搬很重的家具"])]
prompts=["","🤔 还有什么? 说一说!","🤔 还有哪些危险的事?"]
for i,(title,cl,items) in enumerate(cols):
    x=0.35+i*3.15
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.25),Inches(2.95),Inches(3.15))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(3)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.25),Inches(2.95),Inches(0.55))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.1,1.33,2.75,0.42,title,sz=14,b=True,c=WHITE)
    tf=tb(s,x+0.25,1.98,2.55,2.3,f"· {items[0]}",sz=14,b=True,c=DARK)
    for it in items[1:]:ap(tf,"",sz=6);ap(tf,f"· {it}",sz=14,b=True,c=DARK)
    if prompts[i]:tb(s,x+0.25,3.95,2.55,0.35,prompts[i],sz=11,c=GRAY)
warn(s,4.6,"小朋友可以帮忙，但首先要注意安全！不确定的事情，先问爸爸妈妈。")
notes(s,"判断活动(5分钟)：老师逐个说场景，学生用手势表示红/黄/绿(放映时建议先出场景，点击后再出现在对应栏)。绿=叠衣服/收玩具/擦桌子/扫地/洗自己的袜子；黄=洗碗/使用洗衣机；红=使用炉灶/使用锋利的刀/搬很重的家具。强调：可以帮忙，但安全第一。")
pn(s,n)

# ============================================================
# 8 PART 2 DIVIDER
# ============================================================
div("Part 2 · 上午","认识各种家务\n🏠 什么是家务  🗂️ 家务分类  🎭 动作猜家务",AMBER,"🧹")
n+=1

# ============================================================
# 9 什么是家务?
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏠 什么是家务?  What Are Chores?")
defbox=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(0.95),Inches(9.0),Inches(1.0))
defbox.fill.solid();defbox.fill.fore_color.rgb=WARM;defbox.line.color.rgb=PLUM;defbox.line.width=Pt(2.5)
tb(s,0.7,1.08,8.6,0.75,"家务，就是为了让家里保持干净、整齐、舒服，\n需要大家一起完成的事情。",sz=16,b=True,c=PLUM,a=PP_ALIGN.CENTER)
tb(s,0.5,2.15,9,0.4,"✋ 判断游戏：下面哪些是家务? (是→拍手 · 不是→摇头)",sz=15,b=True,c=AMBER)
items=[("🍽️","洗碗",True),("📺","看电视",False),("🧹","扫地",True),
       ("📚","整理书桌",True),("🎮","玩游戏",False),("👕","叠衣服",True)]
for i,(em,txt,yes) in enumerate(items):
    col=i%3;row=i//3
    x=0.55+col*3.15;y=2.7+row*1.15
    cl=GREEN_OK if yes else GRAY
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.0))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    tb(s,x+0.15,y+0.22,0.7,0.55,em,sz=26)
    tb(s,x+0.95,y+0.13,1.9,0.4,txt,sz=16,b=True,c=DARK)
    tb(s,x+0.95,y+0.55,1.9,0.35,"✅ 是家务" if yes else "❌ 不是家务",sz=12,b=True,c=cl)
notes(s,"什么是家务(5分钟)：用儿童语言解释定义。判断游戏：是家务拍手，不是摇头(放映时建议：先出词语，点击后再出现✅/❌)。是：洗碗/扫地/整理书桌/叠衣服；不是：看电视/玩游戏。")
pn(s,n)

# ============================================================
# 10 家务分类
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🗂️ 家务分类  Sort the Chores",PLUM)
tb(s,0.4,0.82,9.2,0.3,"这些家务应该放进哪个「家」? 大声说出来！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
cats=[("🍳","厨房家务","洗碗 · 倒垃圾",AMBER),
      ("👕","衣物整理","叠衣服 · 洗衣服",PLUM),
      ("🧸","房间整理","整理书桌 · 收玩具",BLUE),
      ("🧹","地面桌面清洁","扫地 · 擦桌子",TEAL)]
for i,(em,cn,d,cl) in enumerate(cats):
    x=0.3+i*2.4
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.2),Inches(2.25),Inches(2.5))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    tb(s,x+0.1,1.35,2.05,0.7,em,sz=36,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.15,2.15,0.7,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.95,2.05,0.6,d,sz=12,c=DARK,a=PP_ALIGN.CENTER)
wb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.9),Inches(9.4),Inches(1.0))
wb.fill.solid();wb.fill.fore_color.rgb=WARM;wb.line.color.rgb=PLUM;wb.line.width=Pt(2)
tb(s,0.5,3.98,9.0,0.35,"🎯 词语卡片 (老师念，学生指出正确的家):",sz=13,b=True,c=PLUM)
tb(s,0.5,4.35,9.2,0.42,"洗碗 · 擦桌子 · 扫地 · 叠衣服 · 洗衣服 · 整理书桌 · 收玩具 · 倒垃圾",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
notes(s,"家务分类(5分钟)：四类——厨房家务/衣物整理/房间整理/地面桌面清洁。老师念8个词语，学生口头分类或上台把词卡贴到正确类别(放映时建议先只出四个类别，答案点击出现)。")
pn(s,n)

# ============================================================
# 11 动作猜家务
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 动作猜家务  Act It Out!",AMBER)
tb(s,0.4,0.82,9.2,0.3,"看动作，猜一猜：他/她在做什么? 猜对了才点击出现词语！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
acts=[("🍽️🫧","洗碗"),("👕🙌","叠衣服"),("🧹💨","扫地"),
      ("🧽➡️","擦桌子"),("📚✏️","整理书桌"),("🫧👖","洗衣服")]
for i,(em,txt) in enumerate(acts):
    col=i%3;row=i//3
    x=0.5+col*3.15;y=1.3+row*1.7
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.5))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=AMBER;card.line.width=Pt(2.5)
    tb(s,x,y+0.15,2.95,0.6,em,sz=28,a=PP_ALIGN.CENTER)
    ans=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.55),Inches(y+0.9),Inches(1.85),Inches(0.5))
    ans.fill.solid();ans.fill.fore_color.rgb=PLUM;ans.line.fill.background()
    tb(s,x+0.55,y+0.97,1.85,0.4,txt,sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.4,4.75,9.2,0.35,"💬 句型：他在洗碗。她在叠衣服。",sz=14,b=True,c=DEEP,a=PP_ALIGN.CENTER)
notes(s,"动作猜家务(5分钟)：老师或学生做动作(不出声)，其他人猜「他/她在做什么?」，猜对后点击出现词语(放映时建议给紫色词条加动画)。词语：洗碗/叠衣服/扫地/擦桌子/整理书桌/洗衣服。练习句型「他在…她在…」。")
pn(s,n)

# ============================================================
# 12 SKILL 1 DIVIDER — 洗碗
# ============================================================
div("技能一 · 洗碗","Wash the Dishes\n🍽️ 准备工具  🫧 七个步骤  ✨ 洗干净了吗?",AMBER,"🍽️")
n+=1

# ============================================================
# 13 洗碗前要准备什么?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"1","洗碗前要准备什么?","Get Ready to Wash",AMBER)
tb(s,0.4,0.95,9.2,0.3,"下面哪些东西是洗碗需要的? 指一指，说一说！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
tools=[("🍽️","脏碗"),("🧴","洗碗液"),("🧽","海绵"),("🚰","水槽"),("🧻","抹布"),("🗄️","沥水架")]
for i,(em,txt) in enumerate(tools):
    x=0.4+i*1.58
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(1.45),Inches(1.75))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=AMBER;card.line.width=Pt(2)
    tb(s,x,1.75,1.45,0.65,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,x,2.6,1.45,0.4,txt,sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.65),Inches(9.4),Inches(1.25))
sb.fill.solid();sb.fill.fore_color.rgb=WHITE;sb.line.color.rgb=RED;sb.line.width=Pt(2.5)
tb(s,0.5,3.76,9.0,0.35,"⚠️ 洗碗安全三条  Safety First:",sz=14,b=True,c=RED)
tf=tb(s,0.7,4.16,8.8,0.65,"🔪 小朋友不要洗锋利的刀　　🥛 玻璃制品请大人帮忙　　💧 地上有水要马上擦干",sz=14,b=True,c=DARK)
notes(s,"洗碗准备(5分钟)：展示6样东西，问哪些是洗碗需要的。安全三条：不洗锋利的刀、玻璃制品请大人帮忙、地上有水马上擦干。")
pn(s,n)

# ============================================================
# 14 洗碗步骤
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"2","洗碗七步  Wash Step by Step",  "7 steps to clean dishes",AMBER)
steps7=[("1️⃣","把剩菜倒进垃圾桶"),("2️⃣","用水把碗冲湿"),("3️⃣","海绵上挤一点洗碗液"),("4️⃣","里里外外认真擦洗"),
        ("5️⃣","用清水冲干净"),("6️⃣","放到沥水架上"),("7️⃣","擦干桌面和水槽边的水")]
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.7),Inches(3.85))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=AMBER;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.7),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=AMBER;lh.line.fill.background()
tb(s,0.45,1.12,4.5,0.4,"🍽️ 洗碗顺序  The Order",sz=14,b=True,c=WHITE)
tf=tb(s,0.5,1.68,4.4,3.1,f"{steps7[0][0]} {steps7[0][1]}",sz=13,b=True,c=DARK)
for num,txt in steps7[1:]:ap(tf,"",sz=9);ap(tf,f"{num} {txt}",sz=13,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.05),Inches(4.5),Inches(3.85))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=PLUM;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.05),Inches(4.5),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=PLUM;rh.line.fill.background()
tb(s,5.35,1.12,4.2,0.4,"🎮 排顺序游戏  Put in Order",sz=14,b=True,c=WHITE)
tf2=tb(s,5.35,1.7,4.2,2.0,"把7张步骤图片打乱，",sz=14,b=True,c=DARK)
ap(tf2,"请小朋友按正确顺序",sz=14,b=True,c=DARK)
ap(tf2,"重新排列！",sz=14,b=True,c=DARK)
ap(tf2,"",sz=8)
ap(tf2,"💡 想一想：为什么要先倒剩菜，",sz=13,c=DEEP)
ap(tf2,"　 最后才擦桌面?",sz=13,c=DEEP)
tb(s,5.35,4.35,4.2,0.4,"📷 每步配连续图片更清楚",sz=11,b=True,c=DEEP)
notes(s,"洗碗步骤(8分钟)：七步——倒剩菜/冲湿/挤洗碗液/里外擦洗/清水冲净/放沥水架/擦干桌面水槽。排顺序活动：打乱步骤图片让学生排列。")
pn(s,n)

# ============================================================
# 15 洗干净了吗?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"3","洗干净了吗?","Is It Really Clean?",AMBER)
bowls=[("🥘","碗上还有油","摸起来滑滑的","❌ 没洗干净",RED),
       ("🫧","碗上还有泡沫","洗碗液没冲掉","❌ 没洗干净",RED),
       ("✨","干净又放好","没油·没泡沫·没食物","✅ 真正洗干净!",GREEN_OK)]
for i,(em,t1,t2,t3,cl) in enumerate(bowls):
    x=0.4+i*3.15
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.15),Inches(2.95),Inches(2.9))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(3)
    tb(s,x,1.35,2.95,0.8,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x,2.25,2.95,0.4,t1,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x,2.7,2.95,0.35,t2,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.45),Inches(3.25),Inches(2.05),Inches(0.5))
    pill.fill.solid();pill.fill.fore_color.rgb=cl;pill.line.fill.background()
    tb(s,x+0.45,3.32,2.05,0.4,t3,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
rule=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.3),Inches(9.4),Inches(0.75))
rule.fill.solid();rule.fill.fore_color.rgb=AMBER;rule.line.fill.background()
tb(s,0.5,4.42,9.0,0.5,"🔑 记住：没有食物、没有油、没有泡沫，才算真正洗干净！",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"洗干净了吗(3分钟)：展示三只碗(有油/有泡沫/干净放好)，让学生判断哪只真正洗干净(放映时建议先不出现❌✅标签，学生判断后点击出现)。总结标准：没食物、没油、没泡沫。")
pn(s,n)

# ============================================================
# 16 SKILL 2 DIVIDER — 洗衣服
# ============================================================
div("技能二 · 洗衣服","Do the Laundry\n🧺 先分类  🫧 洗衣机七步  🧦 旅行手洗",BLUE,"🧺")
n+=1

# ============================================================
# 17 洗衣服前先看什么?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"1","洗衣服前先看什么?","Sort Before You Wash",BLUE)
tb(s,0.4,0.95,9.2,0.3,"这些衣服能一起洗吗? 先分一分！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
clothes=[("👕","白衣服"),("👖","深色衣服"),("🧣","毛巾"),("🧦","袜子"),("🧶","毛衣"),("🏷️","特殊标签衣服")]
for i,(em,txt) in enumerate(clothes):
    x=0.4+i*1.58
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(1.45),Inches(1.35))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=BLUE;card.line.width=Pt(2)
    tb(s,x,1.45,1.45,0.6,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x,2.1,1.45,0.55,txt,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
bins=[("⬜","浅色一堆","白衣服 · 浅色T恤",TEAL),
      ("⬛","深色一堆","深色衣服 · 牛仔裤",DEEP),
      ("🙋","请大人帮忙","毛衣 · 有特殊标签的",CORAL)]
for i,(em,cn,d,cl) in enumerate(bins):
    x=0.55+i*3.15
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.95),Inches(2.95),Inches(1.5))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.95),Inches(2.95),Inches(0.5))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.1,3.02,2.75,0.4,f"{em} {cn}",sz=14,b=True,c=WHITE)
    tb(s,x+0.15,3.6,2.65,0.7,d,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,4.65,9.2,0.4,"🔑 洗衣服前要先分类 — 浅色和深色分开洗，颜色才不会染花！",sz=14,b=True,c=BLUE,a=PP_ALIGN.CENTER)
notes(s,"洗衣分类(5分钟)：展示白衣服/深色衣服/毛巾/袜子/毛衣/特殊标签衣服。简单分三堆：浅色/深色/请大人帮忙(毛衣、特殊标签)。不讲复杂洗涤知识，只强调「先分类」。")
pn(s,n)

# ============================================================
# 18 洗衣机洗衣服的步骤
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"2","洗衣机洗衣服七步","Use the Washing Machine",BLUE)
msteps=[("1️⃣","检查口袋","口袋里不能有东西",PLUM),
        ("2️⃣","简单分类","浅色和深色分开",TEAL),
        ("3️⃣","放进洗衣机","不要塞得太满",BLUE),
        ("4️⃣","加洗衣液","适量就好，不要太多",AMBER),
        ("5️⃣","请大人选程序","按钮请大人来按",CORAL),
        ("6️⃣","及时拿出来","洗完别忘在里面",GREEN_OK),
        ("7️⃣","晾干或烘干","挂起来或进烘干机",GOLD)]
for i,(num,cn,d,cl) in enumerate(msteps):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.05+row*0.88
    if i==6:x=2.725  # 最后一个居中
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.68))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2)
    tb(s,x+0.12,y+0.13,0.5,0.45,num,sz=17,b=True)
    tb(s,x+0.7,y+0.14,1.9,0.4,cn,sz=15,b=True,c=cl)
    tb(s,x+2.55,y+0.19,1.95,0.4,d,sz=10,c=DARK)
warn(s,4.6,"小朋友使用洗衣机时，一定要有大人陪同！")
notes(s,"洗衣机七步(8分钟)：检查口袋/简单分类/放进洗衣机/加洗衣液/请大人选程序/及时拿出/晾干或烘干。强调：使用洗衣机必须有大人陪同。")
pn(s,n)

# ============================================================
# 19 为什么要检查口袋?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"3","为什么要检查口袋?","Check the Pockets!",BLUE)
tb(s,0.4,0.95,9.2,0.32,"口袋里的东西跟着衣服一起洗，会发生什么? 猜一猜！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
pockets=[("🧻","纸巾","碎成小片，粘满衣服!"),("🪙","硬币","叮叮当当，伤洗衣机!"),
         ("🖍️","蜡笔","化开染色，衣服全花!"),("🍬","糖果","黏黏糊糊，粘在一起!"),
         ("🧸","小玩具","可能坏掉，堵住机器!"),("✏️","铅笔","可能折断，戳破衣服!")]
for i,(em,txt,bad) in enumerate(pockets):
    col=i%3;row=i//3
    x=0.5+col*3.15;y=1.4+row*1.55
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.95),Inches(1.4))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=BLUE;card.line.width=Pt(2)
    tb(s,x+0.15,y+0.15,0.7,0.55,em,sz=26)
    tb(s,x+0.9,y+0.13,1.95,0.4,txt,sz=15,b=True,c=BLUE)
    tb(s,x+0.15,y+0.75,2.65,0.55,bad,sz=11,c=DARK)
concl=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.6),Inches(9.4),Inches(0.6))
concl.fill.solid();concl.fill.fore_color.rgb=BLUE;concl.line.fill.background()
tb(s,0.5,4.68,9.0,0.45,"🔑 所以：洗衣服前，要先检查每一个口袋！",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"检查口袋(5分钟)：先出示物品(纸巾/硬币/蜡笔/糖果/玩具/铅笔)，让学生猜跟着洗会怎样(放映时建议后果文字点击出现)。总结：洗前检查每个口袋。")
pn(s,n)

# ============================================================
# 20 旅行时手洗小件衣服
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"4","旅行时怎样手洗小件衣服?","Hand-Wash Small Items",BLUE)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(3.2),Inches(3.4))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=TEAL;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(3.2),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=TEAL;lh.line.fill.background()
tb(s,0.45,1.12,3.0,0.4,"🧦 适合手洗的",sz=14,b=True,c=WHITE)
tf=tb(s,0.5,1.9,2.9,2.4,"🧦 袜子",sz=17,b=True,c=DARK)
for it in ["🤧 手帕","🧣 小毛巾"]:
    ap(tf,"",sz=22);ap(tf,it,sz=17,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.7),Inches(1.05),Inches(6.0),Inches(3.4))
rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=BLUE;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.7),Inches(1.05),Inches(6.0),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=BLUE;rh.line.fill.background()
tb(s,3.85,1.12,5.8,0.4,"🫧 手洗七步  7 Steps",sz=14,b=True,c=WHITE)
hsteps=["1️⃣ 接一盆温水","2️⃣ 加一点洗衣液","3️⃣ 把衣物浸湿","4️⃣ 轻轻搓洗脏的地方","5️⃣ 用清水冲干净","6️⃣ 轻轻挤出水","7️⃣ 挂起来晾干"]
tf2=tb(s,3.9,1.68,2.9,2.7,hsteps[0],sz=13,b=True,c=DARK)
for st in hsteps[1:4]:ap(tf2,"",sz=6);ap(tf2,st,sz=13,b=True,c=DARK)
tf3=tb(s,6.85,1.68,2.75,2.7,hsteps[4],sz=13,b=True,c=DARK)
for st in hsteps[5:]:ap(tf3,"",sz=6);ap(tf3,st,sz=13,b=True,c=DARK)
tb(s,3.9,3.85,5.7,0.4,"👐 动作口令：浸一浸 · 搓一搓 · 冲一冲 · 挤一挤 · 挂起来！",sz=12,b=True,c=TEAL)
warn(s,4.6,"不要把洗衣液放进嘴里，洗完衣服要洗手！")
notes(s,"旅行手洗(8分钟)：适合手洗的是袜子/手帕/小毛巾。七步：温水→加洗衣液→浸湿→轻搓脏处→清水冲净→挤水→挂晾。动作口令帮助记忆。安全：洗衣液不入口，洗完洗手。")
pn(s,n)

# ============================================================
# 21 SKILL 3 DIVIDER — 叠衣服和收纳
# ============================================================
div("技能三 · 叠衣服和收纳","Fold & Put Away\n👕 竖式叠衣法  🧍 让T恤站起来  📦 整理抽屉",PLUM,"👕")
n+=1

# ============================================================
# 22 为什么要叠衣服?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"1","为什么要叠衣服?","Why Fold Clothes?",PLUM)
boxes=[("😖","乱塞的抽屉",RED,["衣服皱巴巴","找一件要翻半天","塞不了几件就满了"]),
       ("😊","叠好的抽屉",GREEN_OK,["每件都整整齐齐","一眼就找到","能放更多衣服"])]
for i,(em,title,cl,lines) in enumerate(boxes):
    x=0.3+i*4.85
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.1),Inches(4.6),Inches(2.5))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(3)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.1),Inches(4.6),Inches(0.55))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.15,1.17,4.3,0.42,f"{em} {title}",sz=17,b=True,c=WHITE)
    ib(s,x+0.2,1.8,1.7,1.6,"📷")
    tf=tb(s,x+2.05,1.95,2.4,1.6,f"· {lines[0]}",sz=13,b=True,c=DARK)
    for l in lines[1:]:ap(tf,"",sz=13);ap(tf,f"· {l}",sz=13,b=True,c=DARK)
qb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.15))
qb.fill.solid();qb.fill.fore_color.rgb=WARM;qb.line.color.rgb=PLUM;qb.line.width=Pt(2)
tb(s,0.5,3.93,9.0,0.35,"❓ 比一比  Compare:",sz=13,b=True,c=PLUM)
tb(s,0.5,4.3,9.2,0.6,"哪个抽屉更容易找到衣服? 哪个可以放更多衣服? 哪个看起来更整齐?",sz=14,b=True,c=DARK)
notes(s,"为什么叠衣服(3分钟)：对比乱塞抽屉vs叠好抽屉的图片。三个问题：更容易找到?放更多?更整齐?引出叠衣服的意义。")
pn(s,n)

# ============================================================
# 23 什么是竖式叠衣法?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"2","什么是竖式叠衣法?","Vertical Folding",PLUM)
defb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(9.4),Inches(0.85))
defb.fill.solid();defb.fill.fore_color.rgb=WARM;defb.line.color.rgb=PLUM;defb.line.width=Pt(2.5)
tb(s,0.5,1.16,9.0,0.65,"把衣服叠成一个整齐的小长方形，竖着放进抽屉，\n而不是一件压在一件上面。",sz=15,b=True,c=PLUM,a=PP_ALIGN.CENTER)
fsteps=[("1️⃣","铺平","T恤正面朝下，平铺在桌面上",TEAL),
        ("2️⃣","两边向中间折","左右两边向中间折，袖子也折进去，变成长方形",BLUE),
        ("3️⃣","从下往上折","从衣服下方向上折，可以折两次或三次",AMBER),
        ("4️⃣","立起来检查","能自己站住，就叠得比较整齐！",GREEN_OK)]
for i,(num,cn,d,cl) in enumerate(fsteps):
    y=2.05+i*0.72
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.62))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2)
    tb(s,0.55,y+0.1,0.5,0.42,num,sz=16,b=True)
    tb(s,1.15,y+0.11,2.5,0.4,cn,sz=15,b=True,c=cl)
    tb(s,3.7,y+0.16,5.8,0.4,d,sz=12,c=DARK)
foot=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.95),Inches(9.2),Inches(0.5))
foot.fill.solid();foot.fill.fore_color.rgb=PLUM;foot.line.fill.background()
tb(s,0.5,5.02,9.0,0.38,"📣 动作口令：铺平 —— 折两边 —— 向上折 —— 站起来！",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"竖式叠衣法(10-12分钟总)：先看视频 https://www.youtube.com/watch?v=_I8djreq0LM 。定义：叠成小长方形竖着放。四步：铺平/两边向中间折/从下往上折/立起来检查。全班一起喊动作口令。")
pn(s,n)

# ============================================================
# 24 竖式收纳的好处 + 对比观察
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"3","竖式收纳有什么好处?","Why Fold Vertically?",PLUM)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(3.85))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=GREEN_OK;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=GREEN_OK;lh.line.fill.background()
tb(s,0.45,1.12,4.4,0.4,"👍 六大好处  Benefits",sz=14,b=True,c=WHITE)
bens=["👀 每件衣服都能看见","🙅 不用翻动上面的衣服","✋ 想拿哪件直接拿","🧘 拿走一件，其他不乱","📦 更好利用抽屉空间","🏠 衣服容易放回原位"]
tf=tb(s,0.5,1.72,4.3,3.1,bens[0],sz=13,b=True,c=DARK)
for b_ in bens[1:]:ap(tf,"",sz=12);ap(tf,b_,sz=13,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(3.85))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=CORAL;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=CORAL;rh.line.fill.background()
tb(s,5.25,1.12,4.4,0.4,"🔬 对比观察  Compare",sz=14,b=True,c=WHITE)
tb(s,5.25,1.68,4.35,0.4,"📚 方法一：一件压一件",sz=13,b=True,c=RED)
tf2=tb(s,5.45,2.05,4.15,0.85,"· 最下面有什么衣服?",sz=11,c=DARK)
ap(tf2,"· 想拿最下面的，要怎么做?",sz=11,c=DARK)
ap(tf2,"· 拿出来后，其他会不会乱?",sz=11,c=DARK)
tb(s,5.25,3.05,4.35,0.4,"🧍 方法二：竖着排列",sz=13,b=True,c=GREEN_OK)
tf3=tb(s,5.45,3.42,4.15,0.85,"· 能一次看到所有衣服吗?",sz=11,c=DARK)
ap(tf3,"· 哪种更容易找到想穿的?",sz=11,c=DARK)
ap(tf3,"· 哪种拿取后更容易保持整齐?",sz=11,c=DARK)
tb(s,0.4,5.0,9.2,0.4,"🔑 竖式叠衣法：省空间 + 容易看见 + 容易拿取 + 容易放回！",sz=13,b=True,c=PLUM,a=PP_ALIGN.CENTER)
notes(s,"竖式收纳好处：能看见/不翻动/直接拿/不变乱/省空间/易放回。对比观察：方法一(叠罗汉)提问3个，方法二(竖排)提问3个。教师总结关键句。")
pn(s,n)

# ============================================================
# 25 学生练习 — 让T恤站起来
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"4","学生练习：让T恤站起来!","Make Your T-shirt Stand!",PLUM)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(3.85))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=AMBER;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=AMBER;lh.line.fill.background()
tb(s,0.45,1.12,4.4,0.4,"1️⃣ 第一轮：跟着老师做",sz=14,b=True,c=WHITE)
r1=["① 铺平","② 折左边","③ 折右边","④ 袖子收进去","⑤ 从下往上折","⑥ 竖起来！"]
tf=tb(s,0.5,1.72,4.3,3.0,r1[0],sz=14,b=True,c=DARK)
for r_ in r1[1:]:ap(tf,"",sz=11);ap(tf,r_,sz=14,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(3.85))
rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=GREEN_OK;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=GREEN_OK;rh.line.fill.background()
tb(s,5.25,1.12,4.4,0.4,"2️⃣ 第二轮：自己完成 + 检查",sz=14,b=True,c=WHITE)
r2=["☐ 衣服铺平了","☐ 两边折得差不多宽","☐ 袖子没有露在外面","☐ 叠成整齐的小长方形","☐ 衣服可以竖着站起来!"]
tf2=tb(s,5.3,1.75,4.3,3.0,r2[0],sz=14,b=True,c=DARK)
for r_ in r2[1:]:ap(tf2,"",sz=14);ap(tf2,r_,sz=14,b=True,c=DARK)
notes(s,"学生练习(10分钟)：每人一件儿童T恤。第一轮老师一步步示范学生同步做(铺平/折左/折右/收袖子/向上折/竖起来)。第二轮学生独立完成，用5项检查表自查。不要求每件叠得完全一样——重点是理解「合适大小+竖着排列」。")
pn(s,n)

# ============================================================
# 26 小挑战 — 整齐放进抽屉
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"5","小挑战：整齐放进抽屉","Drawer Challenge",PLUM)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(3.5))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=PLUM;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=PLUM;lh.line.fill.background()
tb(s,0.45,1.12,4.4,0.4,"✅ 摆放要求  The Rules",sz=14,b=True,c=WHITE)
reqs=["➡️ 衣服朝同一个方向","👀 每件衣服都能看见","🤏 之间不要挤得太紧","🧍 拿出一件，其他不倒下"]
tf=tb(s,0.5,1.8,4.3,2.6,reqs[0],sz=14,b=True,c=DARK)
for r_ in reqs[1:]:ap(tf,"",sz=15);ap(tf,r_,sz=14,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(3.5))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=AMBER;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=AMBER;rh.line.fill.background()
tb(s,5.25,1.12,4.4,0.4,"🧠 挑战问题  Think Harder",sz=14,b=True,c=WHITE)
qs=["· 怎样摆放能看到每件衣服的颜色?","· 衣服太厚，折叠次数怎么调整?","· 抽屉很深或很浅，叠法要变吗?","· 哪种衣服比较容易站起来?"]
tf2=tb(s,5.3,1.8,4.3,2.6,qs[0],sz=13,c=DARK)
for q in qs[1:]:ap(tf2,"",sz=16);ap(tf2,q,sz=13,c=DARK)
tb(s,0.4,4.7,9.2,0.4,"💡 提醒：不用叠得完全一样 — 合适的大小 + 竖着排列，就很棒了！",sz=13,b=True,c=DEEP,a=PP_ALIGN.CENTER)
notes(s,"小挑战(5分钟)：把叠好的T恤竖着放进塑料抽屉或收纳盒。四条要求+四个挑战问题。教师提醒：不要求完全一样，重点是理解竖式收纳原理。")
pn(s,n)

# ============================================================
# 27 怎样整理抽屉?
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"6","怎样整理抽屉?","Organize the Drawer",PLUM)
tb(s,0.4,0.95,9.2,0.32,"🔑 同一类物品放在一起，每类有自己的家！",sz=15,b=True,c=PLUM,a=PP_ALIGN.CENTER)
drawers=[("👕","上衣"),("👖","裤子"),("🧦","袜子"),("🩲","内衣"),("🧣","毛巾")]
for i,(em,txt) in enumerate(drawers):
    x=0.55+i*1.85
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.5),Inches(1.7),Inches(1.7))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=PLUM;card.line.width=Pt(2.5)
    tb(s,x,1.68,1.7,0.7,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x,2.5,1.7,0.4,txt,sz=16,b=True,c=PLUM,a=PP_ALIGN.CENTER)
ch=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.55),Inches(9.4),Inches(1.35))
ch.fill.solid();ch.fill.fore_color.rgb=WARM;ch.line.color.rgb=CORAL;ch.line.width=Pt(2.5)
tb(s,0.5,3.65,9.0,0.35,"🎮 课堂挑战  Class Challenge:",sz=14,b=True,c=CORAL)
tb(s,0.5,4.05,9.0,0.4,"这个抽屉乱七八糟 — 上衣、袜子、毛巾全混在一起！",sz=14,b=True,c=DARK)
tb(s,0.5,4.45,9.0,0.4,"请你帮它重新分类，让每类物品回到自己的家。",sz=14,b=True,c=DARK)
notes(s,"整理抽屉(5分钟)：五类——上衣/裤子/袜子/内衣/毛巾。原则：同类放一起。挑战：给一个混乱抽屉(或图片)，学生说出/动手重新分类。")
pn(s,n)

# ============================================================
# 28 SKILL 4 DIVIDER — 擦桌子和扫地
# ============================================================
div("技能四 · 擦桌子和扫地","Wipe & Sweep\n🧽 擦桌子六步  🧹 扫地六步  🧠 任务顺序挑战",TEAL,"🧹")
n+=1

# ============================================================
# 29 擦桌子的步骤
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"1","擦桌子六步","Wipe the Table",TEAL)
wsteps=[("1️⃣","把桌面上的物品拿开"),("2️⃣","把垃圾丢进垃圾桶"),("3️⃣","湿抹布从一边擦到另一边"),
        ("4️⃣","擦桌角和容易漏掉的地方"),("5️⃣","把桌面擦干"),("6️⃣","把物品整齐放回去")]
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.7),Inches(3.85))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=TEAL;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.7),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=TEAL;lh.line.fill.background()
tb(s,0.45,1.12,4.5,0.4,"🧽 擦桌顺序  The Order",sz=14,b=True,c=WHITE)
tf=tb(s,0.5,1.72,4.5,3.0,f"{wsteps[0][0]} {wsteps[0][1]}",sz=13,b=True,c=DARK)
for num,txt in wsteps[1:]:ap(tf,"",sz=12);ap(tf,f"{num} {txt}",sz=13,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.05),Inches(4.5),Inches(3.85))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=CORAL;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.05),Inches(4.5),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=CORAL;rh.line.fill.background()
tb(s,5.35,1.12,4.2,0.4,"👀 观察比较  Which is Better?",sz=14,b=True,c=WHITE)
tb(s,5.35,1.7,4.2,0.4,"方法A：来回乱擦 🌀",sz=14,b=True,c=RED)
tb(s,5.55,2.08,4.0,0.6,"擦过的地方又弄脏了，越擦越乱！",sz=12,c=DARK)
tb(s,5.35,2.75,4.2,0.4,"方法B：从一边到另一边 ➡️",sz=14,b=True,c=GREEN_OK)
tb(s,5.55,3.13,4.0,0.6,"一条一条有顺序，不漏掉也不重复！",sz=12,c=DARK)
tb(s,5.35,3.9,4.2,0.8,"🔑 擦桌子要有方向，\n　 不要来回乱擦。",sz=13,b=True,c=TEAL)
notes(s,"擦桌子(6分钟)：六步——拿开物品/丢垃圾/湿抹布从一边到另一边/擦桌角/擦干/物品放回。观察比较：来回乱擦 vs 有方向地擦，让学生说哪种更好。")
pn(s,n)

# ============================================================
# 30 扫地的步骤
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"2","扫地六步","Sweep the Floor",TEAL)
ssteps=[("1","先收起地上的物品",PLUM),("2","从房间里面往门口扫",TEAL),("3","把垃圾扫成一小堆",BLUE),
        ("4","用簸箕收起来",AMBER),("5","倒进垃圾桶",CORAL),("6","扫把簸箕放回原位",GREEN_OK)]
for i,(num,txt,cl) in enumerate(ssteps):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.1+row*1.05
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.9))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.18),Inches(y+0.19),Inches(0.52),Inches(0.52))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,x+0.18,y+0.22,0.52,0.44,num,sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.9,y+0.24,3.5,0.45,txt,sz=15,b=True,c=DARK)
warn(s,4.45,"扫把是用来扫地的 — 不可以拿扫把追打别人！")
notes(s,"扫地(6分钟)：六步——收物品/从里往门口扫/扫成一小堆/簸箕收/倒垃圾桶/工具放回原位。强调「从里往外」的方向。安全：不拿扫把打闹。")
pn(s,n)

# ============================================================
# 31 任务顺序挑战
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"3","任务顺序挑战","What Comes First?",TEAL)
sc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(9.4),Inches(0.85))
sc.fill.solid();sc.fill.fore_color.rgb=WARM;sc.line.color.rgb=TEAL;sc.line.width=Pt(2.5)
tb(s,0.5,1.15,9.0,0.65,"🏠 情境：地上有玩具，桌上有纸屑，椅子没有放好。\n应该先做什么，再做什么?",sz=14,b=True,c=DEEP,a=PP_ALIGN.CENTER)
order=[("1","收玩具",PLUM),("2","整理桌面",BLUE),("3","擦桌子",AMBER),("4","扫地",CORAL),("5","工具放回去",GREEN_OK)]
for i,(num,txt,cl) in enumerate(order):
    x=0.35+i*1.95
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.15),Inches(1.75),Inches(1.5))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.58),Inches(2.3),Inches(0.6),Inches(0.6))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,x+0.58,2.38,0.6,0.48,num,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x,3.05,1.75,0.4,txt,sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    if i<4:tb(s,x+1.68,2.55,0.35,0.5,"→",sz=20,b=True,c=GRAY)
qb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.95),Inches(9.4),Inches(1.05))
qb.fill.solid();qb.fill.fore_color.rgb=TEAL;qb.line.fill.background()
tb(s,0.5,4.05,9.0,0.4,"💬 用句型说一说:",sz=14,b=True,c=WHITE)
tb(s,0.5,4.45,9.0,0.45,"我先________，再________，最后________。",sz=17,b=True,c=SUNNY,a=PP_ALIGN.CENTER)
notes(s,"任务顺序挑战(5分钟)：情境——地上玩具+桌上纸屑+椅子乱。先让学生讨论顺序(放映时建议参考顺序点击出现)：收玩具→整理桌面→擦桌子→扫地→工具放回。想一想：为什么要先整理再擦扫?(先收拾干净桌面地面才能擦扫)。练习句型「我先…再…最后…」。")
pn(s,n)

# ============================================================
# 32 AFTERNOON DIVIDER — 语言 + 技能站
# ============================================================
div("下午","语言学习 + 技能轮换站 + 家务打卡\n👀 我会认 5 词  📝 我会写 4 字  🏅 四站挑战",DEEP,"📖")
n+=1

# ============================================================
# 33-37 我会认 word cards
# ============================================================
def word_card_read(w,py,en,sent,img):
    global n
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",AMBER)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=PLUM,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=AMBER,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=AMBER;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=AMBER)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    n+=1;pn(s,n);return s
read_words=[
    ("洗","xǐ","wash","我会洗碗，也会洗袜子。","📷 洗"),
    ("叠","dié","fold","我会叠衣服。","📷 叠"),
    ("碗","wǎn","bowl","碗洗得真干净！","📷 碗"),
    ("衣服","yī fu","clothes","叠好的衣服站起来了！","📷 衣服"),
    ("扫地","sǎo dì","sweep","我帮爸爸妈妈扫地。","📷 扫地"),
]
for w,py,en,sent,img in read_words:
    word_card_read(w,py,en,sent,img)

# ============================================================
# 38 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",AMBER)
games=[
    ("1️⃣","看图认词\nPicture Match","出示家务图片\n学生找出对应\n的汉字卡！",WARM),
    ("2️⃣","动作猜词\nCharades","老师做动作\n学生举起正确\n的词卡",WARM),
    ("3️⃣","拍苍蝇\nFly Swatter","字卡贴白板上\n老师说词语\n学生冲上去拍！",PINKBG),
    ("4️⃣","传话筒\nPass the Mic","传球，音乐停\n拿球的人读字卡\n并造一个句子",RGBColor(0xE3,0xF2,0xFD)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
notes(s,"字词游戏(5分钟)：选一个玩。词卡：洗/叠/碗/衣服/扫地。看图认词和动作猜词是本课重点游戏。")
pn(s,n)

# ============================================================
# 39-40 我会写 cards (两字一页)
# ============================================================
def write_pair(chars):
    """chars: list of (字, pinyin, english, 字形提示)"""
    global n
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",PLUM)
    for i,(w,py,en,hint) in enumerate(chars):
        x=0.4+i*4.75
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.0),Inches(4.5),Inches(2.1))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PLUM;sh.line.width=Pt(3)
        tb(s,x+0.15,1.1,2.0,1.5,w,sz=80,b=True,c=PLUM,a=PP_ALIGN.CENTER)
        tb(s,x+2.2,1.25,2.2,0.4,f"{py}",sz=22,b=True,c=AMBER)
        tb(s,x+2.2,1.72,2.2,0.35,en,sz=15,c=GRAY)
        tb(s,x+2.2,2.15,2.25,0.8,f"💡 {hint}",sz=12,b=True,c=DEEP)
        sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.25),Inches(4.5),Inches(1.05))
        sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
        tb(s,x+0.15,3.32,4.2,0.35,"📝 笔顺 Stroke Order",sz=13,b=True,c=PLUM)
        ib(s,x+0.15,3.7,4.2,0.5,"📷 插入笔顺图")
    pr=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.5),Inches(9.2),Inches(0.6))
    pr.fill.solid();pr.fill.fore_color.rgb=PLUM;pr.line.fill.background()
    tb(s,0.6,4.6,9.0,0.4,"✍️ 练习：1️⃣ 空中写 Air Write → 2️⃣ 手心写 Palm Write → 3️⃣ 描红 + 纸上写 3 遍",sz=13,b=True,c=WHITE)
    n+=1;pn(s,n);return s
s_=write_pair([("洗","xǐ","wash","氵三点水 + 先：\n洗东西要用水！"),
               ("衣","yī","clothes","像一件小衣服：\n上面是领子，下面是衣摆")])
notes(s_,"我会写(8分钟)：洗(三点水+先，和水有关)、衣(象形字，像一件衣服)。练习：空中写→手心写→描红+纸上写3遍。每页最多两个汉字。")
s_=write_pair([("扫","sǎo","sweep","扌提手旁 + 彐：\n用手拿扫把扫地！"),
               ("地","dì","floor / ground","土字旁 + 也：\n地面和泥土有关")])
notes(s_,"我会写(8分钟)：扫(提手旁，用手拿扫把)、地(土字旁，和土地有关)。练习：空中写→手心写→描红+纸上写3遍。")

# ============================================================
# 41 我会说 sentence frames
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💬 我会说  Sentence Frames",AMBER)
tb(s,0.4,0.8,9.2,0.3,"看图说句子 — 每人选一句，大声说！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
frames=[("我会帮爸爸妈妈________。","I help my parents ___."),
        ("我可以自己________。","I can ___ by myself."),
        ("做完________以后，我要________。","After ___, I will ___."),
        ("我先________，再________。","First I ___, then I ___."),
        ("今天我负责________。","Today I'm in charge of ___."),
        ("我需要用________来完成这项家务。","I need ___ to do this chore.")]
for i,(cn,en) in enumerate(frames):
    col=i%2;row=i//2
    x=0.35+col*4.75;y=1.25+row*1.22
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.6),Inches(1.08))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=AMBER;card.line.width=Pt(2)
    tb(s,x+0.18,y+0.16,4.25,0.5,cn,sz=15,b=True,c=PLUM)
    tb(s,x+0.18,y+0.68,4.25,0.3,en,sz=11,c=GRAY)
tb(s,0.4,5.02,9.2,0.35,"💡 配家务图片使用：出示图片 → 学生用句型说一句话。",sz=12,b=True,c=DEEP,a=PP_ALIGN.CENTER)
notes(s,"句型练习(8分钟)：六个句型。看图说句子：出示家务图片，学生选句型说完整的话。例：看到扫地图片→「我会帮爸爸妈妈扫地」「我需要用扫把来完成这项家务」。")
pn(s,n)

# ============================================================
# 42 技能轮换站说明
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏅 小小家庭管家技能挑战  Skill Stations",DEEP)
tb(s,0.4,0.82,9.2,0.3,"完成四个技能站打卡，成为小小家庭管家，获得勋章！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
stations=[("1️⃣","衣服收纳站","叠衣服 + 放进抽屉",STA1),
          ("2️⃣","书桌整理站","收拾混乱的书桌",STA2),
          ("3️⃣","厨房清洁站","洗碗 + 擦桌 + 扫地",STA3),
          ("4️⃣","洗衣练习站","手洗小件衣服",STA4)]
for i,(num,cn,d,cl) in enumerate(stations):
    x=0.3+i*2.4
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.2),Inches(2.25),Inches(2.0))
    card.fill.solid();card.fill.fore_color.rgb=cl;card.line.fill.background()
    tb(s,x,1.35,2.25,0.5,num,sz=24,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,1.9,2.15,0.45,cn,sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.45,2.05,0.6,d,sz=11,c=WHITE,a=PP_ALIGN.CENTER)
rules=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.5),Inches(9.4),Inches(1.35))
rules.fill.solid();rules.fill.fore_color.rgb=WHITE;rules.line.color.rgb=DEEP;rules.line.width=Pt(2.5)
tb(s,0.5,3.62,9.0,0.35,"📜 规则  Rules:",sz=14,b=True,c=DEEP)
tf=tb(s,0.7,4.0,8.8,0.85,"⏱️ 每站 10 分钟 · 🔔 听到提示音后轮换 · ✋ 每个人都要动手",sz=14,b=True,c=DARK)
ap(tf,"🏠 使用完工具要放回原位 · 🏅 完成后获得一枚技能印章！",sz=14,b=True,c=DARK)
notes(s,"技能轮换站说明(3分钟)：四站——衣服收纳/书桌整理/厨房清洁/洗衣练习。规则：每站10分钟、提示音轮换、人人动手、工具放回原位、完成得印章。分组：每组4-6人。")
pn(s,n)

# ============================================================
# 43-46 station cards
# ============================================================
def station_card(num,cn,en,color,mats,tasks,extra_title,extra,extra_cl):
    global n
    s=ns();bg(s,CREAM)
    bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.15),Inches(9.4),Inches(0.72))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.45),Inches(0.24),Inches(0.54),Inches(0.54))
    circ.fill.solid();circ.fill.fore_color.rgb=WHITE;circ.line.fill.background()
    tb(s,0.45,0.28,0.54,0.46,num,sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,1.2,0.19,6.5,0.42,cn,sz=21,b=True,c=WHITE)
    tb(s,1.22,0.60,6.5,0.24,en,sz=10,c=WHITE)
    pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.9),Inches(0.27),Inches(1.7),Inches(0.45))
    pill.fill.solid();pill.fill.fore_color.rgb=SUNNY;pill.line.color.rgb=WHITE;pill.line.width=Pt(1.5)
    tb(s,7.95,0.33,1.6,0.36,"⏱️ 10 分钟",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.3),Inches(3.05))
    lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=color;lp.line.width=Pt(2.5)
    lh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.3),Inches(0.5))
    lh2.fill.solid();lh2.fill.fore_color.rgb=color;lh2.line.fill.background()
    tb(s,0.45,1.12,4.1,0.4,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
    tf=tb(s,0.5,1.85,4.0,2.1,mats[0],sz=14,b=True,c=DARK)
    for m in mats[1:]:ap(tf,"",sz=18);ap(tf,m,sz=14,b=True,c=DARK)
    rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.8),Inches(1.05),Inches(4.9),Inches(3.05))
    rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=color;rp.line.width=Pt(2.5)
    rh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.8),Inches(1.05),Inches(4.9),Inches(0.5))
    rh2.fill.solid();rh2.fill.fore_color.rgb=color;rh2.line.fill.background()
    tb(s,4.95,1.12,4.6,0.4,"🎯 任务  Tasks",sz=14,b=True,c=WHITE)
    tf2=tb(s,5.0,1.85,4.6,2.1,tasks[0],sz=14,b=True,c=DARK)
    for t in tasks[1:]:ap(tf2,"",sz=18);ap(tf2,t,sz=14,b=True,c=DARK)
    eb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.35),Inches(9.4),Inches(0.78))
    eb.fill.solid();eb.fill.fore_color.rgb=extra_cl;eb.line.fill.background()
    tb(s,0.5,4.52,9.0,0.5,f"{extra_title} {extra}",sz=13,b=True,c=WHITE)
    n+=1;pn(s,n);return s

s_=station_card("1","第一站 · 衣服收纳站","Clothes Folding Station",STA1,
    ["🗄️ 塑料抽屉","👕 T恤 · 👖 裤子","🧦 袜子 · 🧣 毛巾"],
    ["① 先分类：同类放一堆","② 再折叠：用竖式叠衣法","③ 最后放进正确的抽屉"],
    "🌟 挑战:","在抽屉里留下整齐的空间 — 拿出一件，其他不倒！",STA1)
notes(s_,"第一站衣服收纳站：材料塑料抽屉+T恤裤子袜子毛巾。任务：分类→折叠(竖式)→放进正确抽屉。挑战：留下整齐空间。")
s_=station_card("2","第二站 · 书桌整理站","Desk Cleanup Station",STA2,
    ["📚 书 · ✏️ 铅笔 · 橡皮","🥤 水杯 · 📄 纸张","🧸 玩具 · 🗑️ 垃圾"],
    ["① 丢垃圾 → ② 文具分类","③ 书本放整齐 → ④ 擦桌面","⑤ 常用物品放容易拿的位置"],
    "🌟 挑战:","整理完让同伴 10 秒内找到铅笔和橡皮！",STA2)
notes(s_,"第二站书桌整理站：材料书/铅笔/橡皮/水杯/纸张/玩具/垃圾。任务：丢垃圾→文具分类→书本整齐→擦桌面→常用物品易拿。")
s_=station_card("3","第三站 · 厨房清洁站","Kitchen Cleaning Station",STA3,
    ["🍽️ 塑料碗盘 · 🧽 海绵","🧴 洗碗液 · 🗄️ 沥水架","🧻 抹布 · 🧹 扫把簸箕"],
    ["① 洗碗 → 冲干净 → 放沥水架","② 擦桌子：从一边到另一边","③ 清理地面：扫成一堆倒掉"],
    "⚠️ 安全:","不要跑，不要玩水，不要拿扫把打闹！地上有水马上擦干。",RED)
notes(s_,"第三站厨房清洁站：塑料碗盘/海绵/洗碗液/沥水架/抹布/扫把簸箕。任务：洗碗冲净放架→擦桌→扫地。安全：不跑不玩水不打闹。")
s_=station_card("4","第四站 · 洗衣练习站","Hand-Wash Station",STA4,
    ["🪣 小水盆 · 🧴 少量洗衣液","🧦 袜子或小毛巾","🪝 晾衣架 · 🧣 毛巾"],
    ["① 浸湿 → ② 加少量洗衣液","③ 轻轻搓洗 → ④ 冲干净","⑤ 挤水 → ⑥ 晾起来"],
    "📍 地点:","门口或容易清理水的区域 — 地上有水随时擦干！",STA4)
notes(s_,"第四站洗衣练习站：地点在门口或易清理水的区域。材料小水盆/洗衣液/袜子小毛巾/晾衣架/毛巾。任务：浸湿→加洗衣液→轻搓→冲净→挤水→晾起来。")

# ============================================================
# 47 技能站反思
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💭 技能站反思  Reflection",DEEP)
tb(s,0.4,0.85,9.2,0.3,"四站都完成了！坐下来想一想，说一说：",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
refs=[("🏅","我完成了哪一个任务?",STA1),("😄","哪一个任务最容易?",STA2),
      ("😅","哪一个任务最难?",STA3),("🌱","我学会了什么?",STA4),
      ("🚀","下次我怎样可以做得更好?",AMBER)]
for i,(em,q,cl) in enumerate(refs):
    y=1.3+i*0.74
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.7),Inches(y),Inches(8.6),Inches(0.62))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.88),Inches(y+0.08),Inches(0.46),Inches(0.46))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,0.88,y+0.12,0.46,0.4,em,sz=15,a=PP_ALIGN.CENTER)
    tb(s,1.55,y+0.12,7.6,0.4,q,sz=15,b=True,c=DARK)
tb(s,0.4,5.05,9.2,0.35,"💬 用完整的句子回答，可以用今天学的句型！",sz=12,b=True,c=DEEP,a=PP_ALIGN.CENTER)
notes(s,"技能站反思(8分钟)：五个问题围圈分享。鼓励用句型：「我完成了…」「最容易的是…」「我学会了…」「下次我要…」。")
pn(s,n)

# ============================================================
# 48 我的家务打卡表
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📋 我的家务打卡表  My Chore Chart",PLUM)
tb(s,0.4,0.82,9.2,0.35,"🏠 我在家也能帮忙！把这张表带回家，做一个真正的家庭管家！",sz=14,b=True,c=AMBER,a=PP_ALIGN.CENTER)
# 表头
hdrs=[("家务任务",3.1),("我完成了",1.6),("家长签字",1.6),("我觉得",2.9)]
x=0.4
for htxt,hw in hdrs:
    cell=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(1.35),Inches(hw),Inches(0.5))
    cell.fill.solid();cell.fill.fore_color.rgb=PLUM;cell.line.color.rgb=WHITE;cell.line.width=Pt(1.5)
    tb(s,x,1.42,hw,0.38,htxt,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    x+=hw
chores=["🛏️ 整理自己的床","👕 叠衣服","🧽 擦桌子","🧹 扫地","🧺 把脏衣服放进洗衣篮"]
for r,ch in enumerate(chores):
    y=1.85+r*0.58
    x=0.4
    for ci,(htxt,hw) in enumerate(hdrs):
        cell=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(hw),Inches(0.58))
        cell.fill.solid();cell.fill.fore_color.rgb=WHITE;cell.line.color.rgb=PLUM;cell.line.width=Pt(1)
        if ci==0:tb(s,x+0.12,y+0.12,hw-0.2,0.4,ch,sz=13,b=True,c=DARK)
        elif ci==1:tb(s,x,y+0.08,hw,0.4,"☐",sz=20,c=DARK,a=PP_ALIGN.CENTER)
        elif ci==2:tb(s,x,y+0.13,hw,0.4,"＿＿＿＿",sz=13,c=DARK,a=PP_ALIGN.CENTER)
        else:tb(s,x,y+0.14,hw,0.4,"😀 容易 · 😐 有点难 · 💪 要练习",sz=11,c=DARK,a=PP_ALIGN.CENTER)
        x+=hw
tb(s,0.4,4.85,9.2,0.4,"🖨️ 这张表会打印出来发给每个小朋友 — 黑白彩色都可以！",sz=12,b=True,c=DEEP,a=PP_ALIGN.CENTER)
notes(s,"家务打卡表(5分钟)：介绍表格——5项家务(整理床/叠衣服/擦桌子/扫地/脏衣入篮)，栏目：家务任务/我完成了/家长签字/我觉得(容易·有点难·需要练习)。打印发给学生带回家。也可从备选中换：收拾玩具或书桌/帮忙摆放餐具/给植物浇水。")
pn(s,n)

# ============================================================
# 49 活动规则
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📜 打卡活动规则  How It Works",AMBER)
rules=[("1️⃣","把打卡表带回家",PLUM),
       ("2️⃣","选择 5 项家务完成",BLUE),
       ("3️⃣","每完成一项，请家长签字",TEAL),
       ("4️⃣","周五把打卡表带回来",CORAL),
       ("5️⃣","完成的同学可以参加抓娃娃活动！",AMBER)]
for i,(num,txt,cl) in enumerate(rules):
    y=1.0+i*0.72
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(y),Inches(8.4),Inches(0.6))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2)
    tb(s,1.0,y+0.1,0.6,0.42,num,sz=17,b=True)
    tb(s,1.7,y+0.11,7.3,0.4,txt,sz=16,b=True,c=DARK)
goal=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.7),Inches(9.2),Inches(0.7))
goal.fill.solid();goal.fill.fore_color.rgb=PLUM;goal.line.fill.background()
tb(s,0.6,4.8,8.8,0.5,"❤️ 真正的目标：学会帮助家人，让家里更整齐、更舒服！",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"活动规则(3分钟)：带回家→选5项→家长签字→周五带回→参加抓娃娃。不要过度强调奖励——真正的目标是学会帮助家人。")
pn(s,n)

# ============================================================
# 50 结课总结
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌟 我成为小小家庭管家啦!  I Did It!",DEEP)
qs=[("🏠","什么是家务?"),("💪","你学会了哪些家务?"),("🍽️","洗碗时要注意什么?"),
    ("👖","洗衣服前为什么要检查口袋?"),("🧹","为什么要先整理，再擦桌子扫地?"),("🏡","回家以后，你最想做哪一项家务?")]
for i,(em,q) in enumerate(qs):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.0+row*1.02
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.88))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=PLUM;card.line.width=Pt(2)
    tb(s,x+0.15,y+0.26,0.6,0.5,em,sz=22)
    tb(s,x+0.8,y+0.24,3.6,0.6,q,sz=13,b=True,c=DARK)
slogan=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.25),Inches(9.2),Inches(0.95))
slogan.fill.solid();slogan.fill.fore_color.rgb=AMBER;slogan.line.fill.background()
tb(s,0.6,4.35,8.8,0.4,"📣 结尾口号  Class Chant:",sz=13,b=True,c=WHITE)
tb(s,0.6,4.72,8.8,0.45,"自己的事情自己做，家里的事情一起做！",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"结课总结(5分钟)：六个复习问题快问快答。全班一起喊结尾口号：「自己的事情自己做，家里的事情一起做！」")
pn(s,n)

# ============================================================
# 51 Day 3 Badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.3,9,0.7,"🎖️ Day 3 家庭管家徽章  Home Helper Badge",sz=24,b=True,c=PLUM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.05),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=PLUM;sh.line.width=Pt(5)
tf=tb(s,3.6,1.5,2.8,0.4,"DAY 3",sz=18,b=True,c=AMBER,a=PP_ALIGN.CENTER)
ap(tf,"🏠",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"小小家庭管家",sz=19,b=True,c=PLUM,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=13,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
ap(tf,"🍽️  👕  🧹  🏅",sz=16,a=PP_ALIGN.CENTER)
sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.3),Inches(4.12),Inches(7.4),Inches(0.6))
sb.fill.solid();sb.fill.fore_color.rgb=SUNNY;sb.line.color.rgb=AMBER;sb.line.width=Pt(2.5)
tb(s,1.3,4.15,7.4,0.55,"⭐  ⭐  ⭐  ⭐  ⭐  ⭐",sz=30,b=True,c=AMBER,a=PP_ALIGN.CENTER)
tb(s,1,4.8,8,0.4,"📣 老师:「自己的事情——」 学生:「自己做!」🎉",sz=15,b=True,c=PLUM,a=PP_ALIGN.CENTER)
tb(s,1,5.15,8,0.3,"洗碗 · 洗衣服 · 叠衣收纳 · 擦桌扫地 · 四站打卡 · 5 个词",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
notes(s,"颁发徽章：完成四站打卡的学生获得「家庭管家」徽章或印章。合影留念。")
pn(s,n)

# ============================================================
# 52 教师材料准备清单
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧰 教师材料准备清单  Teacher Prep List",DEEP)
mats=[("🗄️","塑料抽屉"),("👕","儿童衣物 袜子 毛巾"),("📚","混乱书桌用的 书本文具"),("🍽️","塑料碗盘"),
      ("🧽","海绵"),("🧴","少量洗碗液"),("🗄️","沥水架"),("🧻","抹布"),
      ("🧹","扫把和簸箕"),("🪣","小水盆"),("🪝","晾衣架"),("🏅","技能站印章 或贴纸"),
      ("⏱️","计时器"),("📋","打印好的 家务打卡表"),("✍️","家长签字说明"),("🎁","抓娃娃奖品 或代币")]
for i,(em,txt) in enumerate(mats):
    col=i%4;row=i//4
    x=0.32+col*2.37;y=1.0+row*1.03
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.22),Inches(0.85))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=DEEP;card.line.width=Pt(1.5)
    tb(s,x+0.1,y+0.2,0.55,0.5,em,sz=20)
    tb(s,x+0.62,y+0.14,1.56,0.62,txt.replace(" ","\n") if " " in txt else txt,sz=11,b=True,c=DARK)
tb(s,0.4,5.05,9.2,0.35,"💡 洗衣练习站请安排在门口或容易清理水的区域，多备毛巾。",sz=12,b=True,c=RED,a=PP_ALIGN.CENTER)
notes(s,"材料清单16项：塑料抽屉/儿童衣物袜子毛巾/书本文具/塑料碗盘/海绵/洗碗液/沥水架/抹布/扫把簸箕/小水盆/晾衣架/印章贴纸/计时器/打卡表/家长签字说明/抓娃娃奖品。")
pn(s,n)

# ============================================================
# 53 Tomorrow preview
# ============================================================
s=ns();n+=1;bg(s,PLUM)
tb(s,0.5,0.9,9,0.8,"🌟 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.5,"Day 4 — 待定 (To Be Continued)",sz=26,b=True,c=SUNNY,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🏠 继续做小小生活家！",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Keep being a little homemaker!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"别忘了：回家完成你的家务打卡表！",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# --- normalize rounded-rect corner radius (avoids header/panel corner slivers) ---
for slide in prs.slides:
    for sh in slide.shapes:
        try:
            if sh.auto_shape_type==MSO_SHAPE.ROUNDED_RECTANGLE:
                mind=min(sh.width,sh.height)/914400.0
                if mind>0:
                    radius=min(0.16667*mind,0.09)
                    sh.adjustments[0]=radius/mind
        except Exception:pass

OUT='/Users/huanli/projects/courseppt/Chinese/小小生活家/day3_home_manager.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
