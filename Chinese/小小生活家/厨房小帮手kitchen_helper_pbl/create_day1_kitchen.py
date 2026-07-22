#!/usr/bin/env python3
"""
厨房小帮手 Kitchen Helper Unit — Day 1: 厨房安全与我的第一份早餐
Structure modeled on 野外生存与探险 Day 1 (wilderness), adapted for the kitchen.
Palette: Fresh Breakfast (fresh green + orange-yellow) — distinct from wilderness pine.
The "6 environments" of the wilderness unit become the "6 kitchen dangers".
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Fresh Breakfast (清新早餐) ---
FRESH = RGBColor(0x4F,0x9D,0x2F)  # primary: fresh leaf green
ORANGE = RGBColor(0xF2,0x93,0x1E) # accent: warm orange (juice)
YOLK = RGBColor(0xF7,0xC5,0x3B)   # egg-yolk yellow
CREAM = RGBColor(0xFD,0xF7,0xEC)  # background cream
TOMATO = RGBColor(0xE0,0x4A,0x3C) # warning red / ALERT
TOAST = RGBColor(0xB5,0x7A,0x3C)  # toast brown
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
WARM = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xEC,0xEC,0xE6)
SKY = RGBColor(0x19,0x76,0xD2)
ALERT = TOMATO
GREEN_OK = RGBColor(0x38,0x8E,0x3C)

# Traffic-light task colors
GREENL = RGBColor(0x43,0xA0,0x47)  # 绿: child can do alone
YELLOWL = RGBColor(0xF5,0xB4,0x1F) # 黄: needs an adult
REDL = RGBColor(0xE0,0x4A,0x3C)    # 红: cannot do alone

# Per-danger colors
KNIFE = RGBColor(0x78,0x86,0x94)   # steel gray
FIRE = RGBColor(0xE5,0x64,0x2E)    # fire orange-red
HOT = RGBColor(0x2E,0x7D,0xB5)     # hot-water blue
ELEC = RGBColor(0x7B,0x5E,0xA7)    # electric purple
GLASS = RGBColor(0x2A,0xA6,0x9A)   # glass teal
WET = RGBColor(0xC9,0x9A,0x2E)     # caution amber

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=FRESH,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n):
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(9.18),Inches(5.28),Inches(0.5),Inches(0.30))
    chip.fill.solid();chip.fill.fore_color.rgb=WHITE;chip.line.color.rgb=LGRAY;chip.line.width=Pt(0.75)
    bx=s.shapes.add_textbox(Inches(9.18),Inches(5.30),Inches(0.5),Inches(0.26));tf=bx.text_frame;tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    r=p.add_run();r.text=str(n);r.font.size=Pt(10);r.font.color.rgb=GRAY;r.font.name='KaiTi'
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    lines=sub.split("\n")
    tf=tb(s,0.4,2.75,9.2,1.6,lines[0],sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    for ln in lines[1:]:ap(tf,ln,sz=20,c=WHITE,a=PP_ALIGN.CENTER)
    return s

n=0

# ============================================================
# 1 COVER — Kitchen Helper badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.25,8,0.7,"Little Kitchen Helper",sz=32,b=True,c=FRESH,a=PP_ALIGN.CENTER)
tb(s,1,0.85,8,0.45,"厨房小帮手",sz=20,c=FRESH,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.5),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=FRESH;sh.line.color.rgb=ORANGE;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.8),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=ORANGE;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.05,2.8,0.4,"DAY 1",sz=16,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
ap(tf,"🍳",sz=50,a=PP_ALIGN.CENTER)
ap(tf,"厨房安全与早餐",sz=17,b=True,c=FRESH,a=PP_ALIGN.CENTER)
ap(tf,"KITCHEN SAFETY & BREAKFAST",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,5.05,8,0.4,"🧑‍🍳 系好围裙，我们开始！Put on your apron, let's cook!",sz=14,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日时间安排  Today's Schedule")
for i,(nm,tm,dc,cl) in enumerate([
    ("Session 1  上午","11:00-11:45","厨房安全 + 红黄绿任务 + 安全大侦探",FRESH),
    ("Session 2  下午","2:00-2:45","复习 + 语言目标 + 三明治步骤 + 摆餐具",ORANGE),
    ("Session 3  下午","3:00-4:30","三明治制作 + 早餐餐垫 + 回家挑战",TOAST)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,4,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.6,3,0.4,tm,sz=15,c=WARM)
    tb(s,4.6,y+0.35,5.0,0.6,dc,sz=14,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.85,9,0.5,"🍳 内容目标  Content:",sz=19,b=True,c=FRESH)
tf=tb(s,0.7,1.32,9,1.4,"1. 认识厨房中的基本卫生与安全规则",sz=14,c=DARK)
ap(tf,"2. 区分可以独立完成、需要大人帮助、不能独自完成的厨房任务",sz=14,c=DARK)
ap(tf,"3. 学会准备一份简单早餐，并完成食材归位和餐后清洁",sz=14,c=DARK)
ap(tf,"4. 体会家人准备饭菜需要计划、劳动和责任",sz=14,c=DARK)
tb(s,0.5,3.05,9,0.5,"🗣️ 语言目标  Language:",sz=19,b=True,c=ORANGE)
tb(s,0.7,3.5,4.6,0.9,"👀 我会认：厨房 安全 洗手\n　　　　　餐具 早餐 危险 整理",sz=14,b=True,c=DARK)
tb(s,5.4,3.5,4.3,0.9,"✍️ 我会写：安全 洗手",sz=14,b=True,c=DARK)
tb(s,0.5,4.6,9,0.5,"🎨 实践目标：完成 Booklet + 三明治制作 + 早餐餐垫 + 回家挑战",sz=14,c=TOAST)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Session 1  上午","厨房安全 + 红黄绿任务 + 安全大侦探\n🔪 刀  🔥 火  ♨️ 热水  🔌 电器  🥃 玻璃  💦 湿地面",FRESH,"🧼")
n+=1

# ============================================================
# 5 HELPER MISSION — narrative intro
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧑‍🍳 你是厨房小帮手!  You're a Kitchen Helper!",FRESH)
hbx=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(9.2),Inches(0.95))
hbx.fill.solid();hbx.fill.fore_color.rgb=WARM;hbx.line.color.rgb=ORANGE;hbx.line.width=Pt(2.5)
tb(s,0.6,1.10,8.8,0.45,"🌅 今天你要学会在厨房里安全地帮忙、做早餐!",sz=22,b=True,c=FRESH,a=PP_ALIGN.CENTER)
tb(s,0.6,1.55,8.8,0.35,"Today you learn to help safely in the kitchen and make breakfast!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
missions=[("🛡️","认识安全","Learn safety",FRESH),
          ("🚦","红黄绿任务","Sort tasks",YELLOWL),
          ("🥪","做三明治","Make a sandwich",ORANGE),
          ("🍽️","摆好餐桌","Set the table",TOAST)]
for i,(em,cn,en,cl) in enumerate(missions):
    x=0.55+i*2.30;y=2.25
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.10),Inches(1.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,y+0.2,1.9,0.7,em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.95,1.9,0.4,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.32,1.9,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.25),Inches(9.4),Inches(0.95))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2)
tb(s,0.5,4.33,1.7,0.4,"💬 我来说",sz=14,b=True,c=ORANGE)
tb(s,2.0,4.30,7.6,0.3,"我是厨房小帮手，我要帮忙 ____ 。",sz=15,b=True,c=DARK)
tb(s,2.0,4.62,7.6,0.3,"I'm a kitchen helper. I can help ___.",sz=11,c=GRAY)
tb(s,2.0,4.88,7.6,0.3,"系好围裙，敬个礼! Put on your apron and salute! 🫡",sz=12,b=True,c=FRESH)
pn(s,n)

# ============================================================
# 6 情境导入 — Who is making breakfast? (Story / scenario hook)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌅 谁在准备早餐?  Who Makes Breakfast?",ORANGE)
ib(s,0.3,1.0,4.4,3.6,"📷 忙碌的早晨 / A busy morning")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.0),Inches(4.8),Inches(3.6))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=FRESH;sh.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.0),Inches(4.8),Inches(0.5))
head.fill.solid();head.fill.fore_color.rgb=FRESH;head.line.fill.background()
tb(s,5.05,1.08,4.6,0.4,"👀 看一看，想一想  Look & Think",sz=15,b=True,c=WHITE)
tf=tb(s,5.1,1.7,4.55,0.4,"🍳 早晨的家里，谁在忙？",sz=14,c=DARK)
ap(tf,"",sz=6)
ap(tf,"🧺 准备早餐要做哪些事情？",sz=14,c=DARK)
ap(tf,"",sz=6)
ap(tf,"⏰ 早晨为什么这么忙？",sz=14,c=DARK)
ap(tf,"",sz=6)
ap(tf,"🤝 你可以帮忙做什么？",sz=14,c=DARK)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.75),Inches(9.4),Inches(0.5))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2)
tb(s,0.5,4.83,9.0,0.4,"📌 准备早餐不只是做食物 — 还要计划、动手、收拾！",sz=14,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
notes(s,"情境导入：出示一张忙碌早晨的图片。让孩子说出图里的人在做什么。引出：做早餐需要很多步骤和帮手。")
pn(s,n)

# ============================================================
# 6a 早餐的完整流程 — 3 zones (before / during / after)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧺 一顿早餐的完整流程  A Full Breakfast Routine",ORANGE)
tb(s,0.4,0.82,9.2,0.32,"除了制作食物，开始前和吃完后还需要完成哪些事情？",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
zones=[
    ("🧺","准备前","Before",FRESH,["想好吃什么","看家里有什么食材","准备食材和工具","洗手","洗水果和蔬菜","摆放碗筷"]),
    ("🍳","制作中","During",ORANGE,["切、搅拌、加热","注意时间","注意安全","边做边整理","把用完的食材放回去"]),
    ("🧽","吃完后","After",TOAST,["收拾碗筷","保存剩余食物","擦桌子和台面","洗碗","清理垃圾","把厨房恢复整洁"]),
]
for i,(em,cn,en,cl,items) in enumerate(zones):
    x=0.3+i*3.15
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.28),Inches(3.0),Inches(3.4))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.28),Inches(3.0),Inches(0.55))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.12,1.34,0.55,0.45,em,sz=22,c=WHITE)
    tb(s,x+0.72,1.37,1.6,0.42,cn,sz=17,b=True,c=WHITE)
    tb(s,x+0.7,1.42,2.15,0.32,en,sz=10,c=WARM,a=PP_ALIGN.RIGHT)
    tf=tb(s,x+0.2,1.95,2.7,0.35,f"· {items[0]}",sz=12,c=DARK)
    for it in items[1:]:ap(tf,f"· {it}",sz=12,c=DARK)
bs=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.85),Inches(9.4),Inches(0.5))
bs.fill.solid();bs.fill.fore_color.rgb=WARM;bs.line.color.rgb=ORANGE;bs.line.width=Pt(2)
tb(s,0.45,4.93,9.0,0.4,"💡 原来做早餐不只是「做饭」— 还要计划、准备、安全、清洁和整理！",sz=13,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
notes(s,"在黑板上放三个区域：准备前 / 制作中 / 吃完后。让孩子把各项工作说出来并归类。")
pn(s,n)

# ============================================================
# 6b 分层活动 — sort (younger) + discuss (older)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧩 分一分，想一想  Sort & Think",FRESH)
GOLD=RGBColor(0xD1,0x8F,0x0A)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.0),Inches(4.4),Inches(3.95))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=GOLD;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.0),Inches(4.4),Inches(0.55))
lh.fill.solid();lh.fill.fore_color.rgb=GOLD;lh.line.fill.background()
tb(s,0.45,1.07,4.2,0.42,"🐣 小孩子  Younger",sz=15,b=True,c=WHITE)
tb(s,0.5,1.72,4.0,0.4,"把图片卡放到正确的区域：",sz=13,b=True,c=DARK)
mini=[("🧺","准备前",FRESH),("🍳","制作中",ORANGE),("🧽","吃完后",TOAST)]
for j,(em,cn,cl) in enumerate(mini):
    y=2.28+j*0.70
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(4.0),Inches(0.56))
    chip.fill.solid();chip.fill.fore_color.rgb=cl;chip.line.fill.background()
    tb(s,0.65,y+0.09,3.8,0.4,f"{em} {cn}",sz=15,b=True,c=WHITE)
tb(s,0.5,4.5,4.0,0.32,"💬 这个是「____」做的。",sz=12,b=True,c=GOLD)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.0),Inches(4.8),Inches(3.95))
rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=FRESH;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.0),Inches(4.8),Inches(0.55))
rh.fill.solid();rh.fill.fore_color.rgb=FRESH;rh.line.fill.background()
tb(s,5.05,1.07,4.5,0.42,"🐔 大孩子  Older — 一起讨论",sz=15,b=True,c=WHITE)
dq=["哪些事情必须先做？","哪些可以同时完成？","哪些最花时间？","哪些看不见，却要提前想到？","如果少做一步，会怎样？","怎样分工最公平、最有效率？"]
tf=tb(s,5.05,1.72,4.55,0.4,f"❓ {dq[0]}",sz=13,c=DARK)
for q in dq[1:]:ap(tf,"",sz=4);ap(tf,f"❓ {q}",sz=13,c=DARK)
notes(s,"小孩子：把图片卡放到三个区域。大孩子：讨论顺序、并行、耗时、看不见的工作、少一步的后果、分工。")
pn(s,n)

# ============================================================
# 6c 体会家人的辛劳 — discover, don't preach
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💗 家人辛苦吗?  Let's Count & Reflect",TOMATO)
cc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.92),Inches(9.4),Inches(0.82))
cc.fill.solid();cc.fill.fore_color.rgb=WARM;cc.line.color.rgb=TOMATO;cc.line.width=Pt(2.5)
tb(s,0.5,0.98,9.0,0.42,"🧮 刚才我们一共找到了 ______ 项工作!",sz=19,b=True,c=TOMATO,a=PP_ALIGN.CENTER)
tb(s,0.5,1.40,9.0,0.3,"把孩子说的工作都写在白板上 — 常常有十几项甚至更多！",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.9),Inches(4.7),Inches(3.1))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=TOMATO;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.9),Inches(4.7),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=TOMATO;lh.line.fill.background()
tb(s,0.45,1.96,4.5,0.4,"🤔 一起想一想  Reflect",sz=14,b=True,c=WHITE)
rq=["原来做早餐只有「做饭」一件事吗？","哪些工作我们平时没注意到？","一个人做完所有事，会怎样？","家人每天准备食物，容易吗？","除了说谢谢，还能怎样感谢？"]
tf=tb(s,0.48,2.5,4.5,0.4,f"· {rq[0]}",sz=12,c=DARK)
for q in rq[1:]:ap(tf,"",sz=3);ap(tf,f"· {q}",sz=12,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.9),Inches(4.7),Inches(3.1))
rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=FRESH;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.9),Inches(4.7),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=FRESH;rh.line.fill.background()
tb(s,5.15,1.96,4.5,0.4,"🐔 大孩子开放题  Older",sz=14,b=True,c=WHITE)
oq=["家务怎样分工才算公平？","「公平」是每人做一样吗？","不同年龄，怎样承担不同责任？","家里有哪些「看不见的工作」？"]
tf=tb(s,5.18,2.48,4.45,0.4,f"· {oq[0]}",sz=12,c=DARK)
for q in oq[1:]:ap(tf,"",sz=2);ap(tf,f"· {q}",sz=12,c=DARK)
iw=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.15),Inches(4.05),Inches(4.4),Inches(0.82))
iw.fill.solid();iw.fill.fore_color.rgb=WARM;iw.line.color.rgb=FRESH;iw.line.width=Pt(1.5)
tb(s,5.28,4.10,4.2,0.72,"👀 看不见的工作：提前计划、买菜、记住家人的饮食、检查食物过期、收好剩菜。",sz=11,b=True,c=DARK)
notes(s,"不要说教。让孩子通过「任务数量」自己发现家人的辛劳。数一数白板上的工作，再问反思问题。")
pn(s,n)

# ============================================================
# 6d 早餐小帮手 · 不重复接力赛 — closing group challenge
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎤 早餐小帮手 · 不重复接力赛  Helper Relay!",FRESH)
rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.92),Inches(9.4),Inches(0.55))
rb.fill.solid();rb.fill.fore_color.rgb=WARM;rb.line.color.rgb=FRESH;rb.line.width=Pt(2)
tb(s,0.45,0.98,9.0,0.42,"🔁 每组轮流说一个自己能做的任务 — 不能和前面小组重复！",sz=14,b=True,c=FRESH,a=PP_ALIGN.CENTER)
tb(s,0.4,1.58,9.2,0.3,"🗣️ 我可以帮忙做的事 (例子):",sz=12,b=True,c=DARK)
ex=["🍓 洗水果","🍽️ 摆碗筷","🧽 擦桌子","🥛 把牛奶放回冰箱","🍴 收拾自己的盘子","💦 提醒地上有水","🗑️ 分类垃圾"]
for j,txt in enumerate(ex):
    col=j%4;row=j//4
    x=0.3+col*2.37;y=1.95+row*0.6
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.25),Inches(0.5))
    chip.fill.solid();chip.fill.fore_color.rgb=WHITE;chip.line.color.rgb=FRESH;chip.line.width=Pt(1.5)
    tb(s,x+0.12,y+0.09,2.05,0.35,txt,sz=12,b=True,c=DARK)
tb(s,0.4,3.18,9.2,0.3,"🏆 计分 — 3 种星星 (不只按说得多):",sz=12,b=True,c=ORANGE)
stars=[("说得具体","specific"),("做法安全","safe"),("说明了原因","gives a reason")]
for j,(cn,en) in enumerate(stars):
    x=0.3+j*3.15
    sc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(3.5),Inches(3.0),Inches(0.72))
    sc.fill.solid();sc.fill.fore_color.rgb=YOLK;sc.line.color.rgb=ORANGE;sc.line.width=Pt(1.5)
    tb(s,x+0.15,3.62,0.5,0.5,"⭐",sz=22)
    tb(s,x+0.72,3.58,2.2,0.35,cn,sz=14,b=True,c=DARK)
    tb(s,x+0.72,3.90,2.2,0.28,en,sz=9,c=TOAST)
oe=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.4),Inches(9.4),Inches(0.85))
oe.fill.solid();oe.fill.fore_color.rgb=WARM;oe.line.color.rgb=ORANGE;oe.line.width=Pt(2)
tb(s,0.5,4.47,3.0,0.35,"🐔 大孩子要扩展 + 说原因:",sz=12,b=True,c=ORANGE)
tb(s,0.5,4.80,9.2,0.4,"「我可以在吃早餐前把水果洗干净，因为这件事安全，而且能帮家人节省时间。」",sz=12,b=True,c=DARK)
notes(s,"结尾接力赛：每组说一个不重复的任务。3 星计分：说得具体 / 做法安全 / 说明原因。大孩子必须扩展并说原因。")
pn(s,n)

# ============================================================
# 7 6 DANGERS OVERVIEW
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⚠️ 厨房里有什么危险?  Kitchen Dangers")
tb(s,0.4,0.9,9,0.4,"厨房很有用，但也有危险 — 我们先认识 6 个！",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
dangers=[
    ("🔪","刀具","Knives",KNIFE),
    ("🔥","炉火","Stove & Fire",FIRE),
    ("♨️","热水","Hot Water",HOT),
    ("🔌","电器","Appliances",ELEC),
    ("🥃","碎玻璃","Broken Glass",GLASS),
    ("💦","湿滑地面","Wet Floor",WET),
]
for i,(em,cn,en,cl) in enumerate(dangers):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.45+row*1.95
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,y+0.12,2.8,0.7,em,sz=38,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.95,2.8,0.45,cn,sz=21,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+1.42,2.8,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 8-13  DANGER CARDS — one per danger
# ============================================================
def danger_card(em,cn,en,color,danger,happens,rule,level,level_color,frame):
    global n
    s=ns();bg(s,CREAM)
    # Header bar with identity + red/yellow/green level pill
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.15),Inches(9.4),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.5,0.22,1.0,0.55,em,sz=28,c=WHITE)
    tb(s,1.5,0.20,3.2,0.5,cn,sz=26,b=True,c=WHITE)
    tb(s,1.5,0.62,3.2,0.25,en,sz=11,c=WARM)
    pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.6),Inches(0.27),Inches(4.0),Inches(0.45))
    pill.fill.solid();pill.fill.fore_color.rgb=level_color;pill.line.color.rgb=WHITE;pill.line.width=Pt(1.5)
    tb(s,5.7,0.32,3.8,0.4,level,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # Image placeholder (left)
    ib(s,0.3,1.05,4.3,3.3,f"📷 {cn} 图片")
    # Info panel (right) — 3 stacked rows
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(3.3))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=color;panel.line.width=Pt(2.5)
    rows=[("⚠️ 危险在哪?",danger,TOMATO),("😣 会怎样?",happens,ORANGE),("✅ 怎么做?",rule,GREEN_OK)]
    for i,(lbl,txt,lc) in enumerate(rows):
        y=1.20+i*1.03
        tb(s,5.05,y,4.5,0.3,lbl,sz=14,b=True,c=lc)
        tb(s,5.05,y+0.33,4.55,0.65,txt,sz=13,c=DARK)
    # Sentence frame at bottom
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.5),Inches(9.4),Inches(0.65))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2)
    tb(s,0.5,4.60,1.7,0.4,"💬 我来说",sz=14,b=True,c=ORANGE)
    tb(s,2.0,4.60,7.6,0.4,frame,sz=14,c=DARK)
    n+=1;pn(s,n)
    return s

danger_card("🔪","刀具","Knives",KNIFE,
    "刀很锋利，刀刃会割伤手指。",
    "手被割破会流血、很疼。",
    "刀是红色任务，只有大人能用；不碰刀架上的刀。",
    "🔴 红色 · 不能独自使用",REDL,
    "刀很危险，因为 ____ 。我应该请大人帮忙。")
danger_card("🔥","炉火","Stove & Fire",FIRE,
    "炉火和炉子很烫，还可能引起火灾。",
    "会烧伤皮肤，衣服也可能着火。",
    "不自己开炉火；离炉子远一点；有火告诉大人。",
    "🔴 红色 · 不能独自使用",REDL,
    "我不开炉火，因为 ____ 。")
danger_card("♨️","热水","Hot Water",HOT,
    "热水、热汤、刚烧好的水很烫。",
    "会烫伤手和嘴，很疼。",
    "不自己端热水；请大人帮忙倒；先摸杯子外面试温度。",
    "🔴 红色 · 请大人帮忙",REDL,
    "热水很危险，我请大人 ____ 。")
danger_card("🔌","电器","Appliances",ELEC,
    "插座、烤箱、微波炉可能触电或很烫。",
    "会触电或烫伤，还可能坏掉。",
    "手干了才碰插头；不乱按电器；不把手伸进烤箱。",
    "🟡 黄色 · 需要大人陪伴",YELLOWL,
    "用电器要 ____ ，还要大人在旁边。")
danger_card("🥃","碎玻璃","Broken Glass",GLASS,
    "玻璃杯、玻璃碗打碎会有尖尖的碎片。",
    "碎片会划伤手和脚。",
    "不自己捡碎片；站着别动，告诉大人来清理。",
    "🔴 红色 · 告诉大人清理",REDL,
    "玻璃碎了，我应该 ____ 。")
danger_card("💦","湿滑地面","Wet Floor",WET,
    "地上有水、有油会很滑。",
    "会滑倒、摔伤。",
    "看到水擦干净或告诉大人；在厨房里走路不跑。",
    "🟢 绿色 · 我能帮忙擦干",GREENL,
    "地上有水，我可以帮忙 ____ 。")

# ============================================================
# 13b 厨房安全规则 — Kitchen Safety Rules (2x2 reference)
#     Source: kids-cooking-activities.com/kitchen-safety-rules.html
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🛡️ 厨房安全规则  Kitchen Safety Rules",FRESH)
tb(s,0.4,0.82,9.2,0.32,"记住这些规则，做一个安全的厨房小帮手！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
rule_cards=[
    ("📋","基本规则","Basic Rules",FRESH,
     ["有大人在才能做饭","在厨房里走，不要跑","用完食材放回原处","太烫、太重就请大人帮忙"]),
    ("🔪","用刀安全","Knife Safety",KNIFE,
     ["刀放在拿不到的地方","只在大人陪同下用刀","手指远离刀刃","刀掉了不要去接"]),
    ("🔥","炉火 · 烤箱","Stove & Oven",FIRE,
     ["锅柄转向里面","用围裙和隔热手套","不越过热锅拿东西","大人先开炉火和烤箱"]),
    ("🧼","洗手 · 卫生","Wash & Hygiene",HOT,
     ["做饭前用肥皂洗手 20 秒","生食和熟食分开","不舔手指和勺子","长头发扎起来"]),
]
for i,(em,cn,en,cl,rules) in enumerate(rule_cards):
    col=i%2;row=i//2
    x=0.3+col*4.8;y=1.24+row*1.86
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.6),Inches(1.78))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.6),Inches(0.42))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.12,y+0.05,0.55,0.35,em,sz=18,c=WHITE)
    tb(s,x+0.7,y+0.06,2.4,0.32,cn,sz=15,b=True,c=WHITE)
    tb(s,x+3.0,y+0.10,1.5,0.28,en,sz=10,c=WARM,a=PP_ALIGN.RIGHT)
    tf=tb(s,x+0.18,y+0.50,4.3,0.3,f"· {rules[0]}",sz=11,c=DARK)
    for r in rules[1:]:ap(tf,f"· {r}",sz=11,c=DARK)
bs=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.02),Inches(9.4),Inches(0.4))
bs.fill.solid();bs.fill.fore_color.rgb=WARM;bs.line.color.rgb=ORANGE;bs.line.width=Pt(1.5)
tb(s,0.45,5.07,9.0,0.32,"💡 把这些规则做成海报，贴在厨房墙上!  Make a poster for the kitchen wall!",sz=12,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 14 厨房安全大侦探 — spot the danger
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔍 厨房安全大侦探  Kitchen Safety Detective",ORANGE)
tb(s,0.4,0.9,9.2,0.35,"看「问题厨房」的图片 — 找出危险，说出正确的做法！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.3,1.35,5.4,3.4,"📷 问题厨房 / A messy, dangerous kitchen")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(1.35),Inches(3.8),Inches(3.4))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=ORANGE;panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(1.35),Inches(3.8),Inches(0.5))
head.fill.solid();head.fill.fore_color.rgb=ORANGE;head.line.fill.background()
tb(s,6.05,1.43,3.5,0.4,"🕵️ 你能找到几个?",sz=14,b=True,c=WHITE)
tf=tb(s,6.05,2.0,3.6,0.4,"🔪 刀放在桌边？",sz=13,c=DARK)
for q in ["🔥 炉火没关？","💦 地上有水？","🥃 玻璃杯要掉了？","🔌 湿手碰插头？"]:
    ap(tf,"",sz=5);ap(tf,q,sz=13,c=DARK)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.9),Inches(9.4),Inches(0.5))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=TOMATO;sf.line.width=Pt(2)
tb(s,0.5,4.97,9.0,0.4,"💬 这里很危险，因为 ____ 。我应该 ____ 。",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
notes(s,"用一张危险厨房图片（可网上找或自制）。让孩子轮流上来指出危险并说正确做法。找到越多越好。")
pn(s,n)

# ============================================================
# 15 红黄绿 CONCEPT
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🚦 厨房任务红黄绿  Red · Yellow · Green Tasks",FRESH)
tb(s,0.4,0.9,9.2,0.35,"厨房里的任务分三种颜色 — 你能做哪些？",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
cols=[("🟢","绿色 Green","我可以独立完成","洗手 · 摆餐具 · 擦桌子 · 倒麦片",GREENL),
      ("🟡","黄色 Yellow","需要大人陪伴","切软水果 · 拿鸡蛋 · 倒牛奶",YELLOWL),
      ("🔴","红色 Red","不能独自完成","开炉火 · 用大刀 · 端热汤",REDL)]
for i,(em,title,mean,ex,cl) in enumerate(cols):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(3.0),Inches(3.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(3.0),Inches(0.6))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.1,1.42,2.8,0.45,f"{em} {title}",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,2.15,2.7,0.7,mean,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER)
    ls=ex.split(" · ")
    tf=tb(s,x+0.2,3.0,2.6,0.4,ls[0],sz=14,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=14,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 16-17 红黄绿 SORTING GAME — Question → Reveal
# ============================================================
SORT_TASKS=[("🧼","洗手","绿"),("🍽️","摆餐具","绿"),("🧽","擦桌子","绿"),
            ("🍓","切软水果","黄"),("🥚","开冰箱拿鸡蛋","黄"),("🥛","倒牛奶","黄"),
            ("🔥","开炉火","红"),("🍲","端热汤","红"),("🔪","用大刀","红")]
CAT_COLOR={"绿":GREENL,"黄":YELLOWL,"红":REDL}
CAT_LABEL={"绿":"🟢 绿","黄":"🟡 黄","红":"🔴 红"}
def sort_game(reveal):
    global n
    s=ns();bg(s,CREAM)
    if reveal:
        hb(s,"🚦 红黄绿分类  💡 答案揭晓 Reveal!",FRESH)
    else:
        hb(s,"🚦 红黄绿分类  🤔 你觉得呢? Think!",FRESH)
    rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(9.4),Inches(0.62))
    rb.fill.solid();rb.fill.fore_color.rgb=WARM;rb.line.color.rgb=ORANGE;rb.line.width=Pt(2)
    tb(s,0.45,1.02,9.0,0.5,"🟢 自己做   🟡 大人陪   🔴 不能自己做   —   老师念任务，学生举手指颜色！",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for i,(em,task,cat) in enumerate(SORT_TASKS):
        col=i%3;row=i//3
        x=0.3+col*3.2;y=1.75+row*1.05
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(0.92))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE
        sh.line.color.rgb=CAT_COLOR[cat] if reveal else GRAY;sh.line.width=Pt(2 if reveal else 1.5)
        tb(s,x+0.12,y+0.24,0.6,0.5,em,sz=24)
        tb(s,x+0.72,y+0.10,1.7,0.72,task,sz=13,b=True,c=DARK)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+2.42),Inches(y+0.26),Inches(0.42),Inches(0.42))
        if reveal:
            badge.fill.solid();badge.fill.fore_color.rgb=CAT_COLOR[cat];badge.line.fill.background()
            tb(s,x+2.42,y+0.30,0.42,0.35,cat,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        else:
            badge.fill.solid();badge.fill.fore_color.rgb=RGBColor(0xEE,0xEE,0xEE);badge.line.color.rgb=GRAY;badge.line.width=Pt(1)
            tb(s,x+2.42,y+0.28,0.42,0.4,"?",sz=18,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    bs=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(5.0),Inches(9.4),Inches(0.42))
    if reveal:
        bs.fill.solid();bs.fill.fore_color.rgb=GREEN_OK;bs.line.fill.background()
        tb(s,0.45,5.05,9.0,0.34,"💡 看看你分对了吗? G1-3 加一句:「这是___色，因为___」",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    else:
        bs.fill.solid();bs.fill.fore_color.rgb=ORANGE;bs.line.fill.background()
        tb(s,0.45,5.05,9.0,0.34,"👉 老师念任务，全班指颜色！🟢自己 🟡大人陪 🔴不能   🤔 3...2...1!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    notes(s,"玩法(3-5分钟):Question页念9个任务，学生用手指或举颜色卡表态；Reveal页揭晓。答错不出局，再听一次。G1-3加因果句。")
    n+=1;pn(s,n);return s
sort_game(False)
sort_game(True)

# ============================================================
# 18 讨论 — responsibility & empathy
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💗 想一想，聊一聊  Let's Discuss",ORANGE)
qs=[("1️⃣","准备早餐，除了做食物，还需要做什么？","Besides cooking, what else is needed?",FRESH),
    ("2️⃣","家长每天准备饭菜，容易吗？为什么？","Is it easy for parents? Why?",TOMATO)]
for i,(num,q_cn,q_en,cl) in enumerate(qs):
    y=1.0+i*1.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.6),Inches(y+0.18),Inches(0.65),Inches(0.65))
    nb.fill.solid();nb.fill.fore_color.rgb=cl;nb.line.fill.background()
    tb(s,0.6,y+0.22,0.65,0.55,num,sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.5,y+0.13,8.0,0.5,q_cn,sz=18,b=True,c=DARK)
    tb(s,1.5,y+0.62,8.0,0.35,q_en,sz=12,c=GRAY)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.4),Inches(9.2),Inches(1.85))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.4),Inches(9.2),Inches(0.45))
head.fill.solid();head.fill.fore_color.rgb=ORANGE;head.line.fill.background()
tb(s,0.6,3.45,9.0,0.4,"💬 我会说  Sentence Frames",sz=14,b=True,c=WHITE)
frames=["准备早餐还要 __________。","家人准备饭菜要 __________。",
        "我觉得 __________ ，因为 __________。","我可以帮忙 __________。"]
for i,fr in enumerate(frames):
    col=i%2;row=i//2
    x=0.6+col*4.55;y=3.95+row*0.55
    tb(s,x,y,4.4,0.4,f"·  {fr}",sz=13,c=DARK)
notes(s,"引导孩子体会家人的劳动与责任。鼓励说出计划、动手、收拾三方面。")
pn(s,n)

# ============================================================
# 19 COMPARISON TABLE — 6 dangers
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🗂️ 6 个危险对比  Compare 6 Dangers",ORANGE)
tb(s,0.4,0.82,9,0.3,"每个危险都不一样，你能记住安全做法吗？",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
ts=s.shapes.add_table(4,7,Inches(0.3),Inches(1.2),Inches(9.4),Inches(3.9));t=ts.table
t.columns[0].width=Inches(1.4)
for i in range(1,7):t.columns[i].width=Inches(1.333)
rows=[
    ["","🔪 刀具","🔥 炉火","♨️ 热水","🔌 电器","🥃 玻璃","💦 湿地"],
    ["😣 会怎样","割伤","烧伤","烫伤","触电\n烫伤","划伤","滑倒"],
    ["🚦 红黄绿","🔴 红","🔴 红","🔴 红","🟡 黄","🔴 红","🟢 绿"],
    ["✅ 怎么做","给大人","不开火","请人倒","手干才碰","叫大人","擦干"],
]
for r,rd in enumerate(rows):
    for c,ct in enumerate(rd):
        cl=t.cell(r,c);cl.text="";tf=cl.text_frame;tf.word_wrap=True
        p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
        rn=p.add_run();rn.text=ct.split('\n')[0];rn.font.name='KaiTi'
        rn.font.size=Pt(13 if r==0 else 11);rn.font.bold=(r==0 or c==0)
        for line in ct.split('\n')[1:]:
            p2=tf.add_paragraph();p2.alignment=PP_ALIGN.CENTER
            rn2=p2.add_run();rn2.text=line;rn2.font.name='KaiTi';rn2.font.size=Pt(11);rn2.font.color.rgb=DARK
        if r==0:
            rn.font.color.rgb=WHITE;cl.fill.solid();cl.fill.fore_color.rgb=FRESH
        elif c==0:
            rn.font.color.rgb=DARK;cl.fill.solid();cl.fill.fore_color.rgb=WARM
        else:
            rn.font.color.rgb=DARK
            if r%2==0:cl.fill.solid();cl.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF0)
pn(s,n)

# ============================================================
# 20-21 安全 vs 危险 gesture game — Question → Reveal
# ============================================================
def gesture_game(reveal):
    global n
    s=ns();bg(s,CREAM)
    if reveal:
        hb(s,"🛡️ 安全 vs 危险!  💡 答案揭晓 Reveal!",TOMATO)
    else:
        hb(s,"🛡️ 安全 vs 危险!  🤔 你觉得呢? Think!",TOMATO)
    rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(9.4),Inches(1.05))
    rb.fill.solid();rb.fill.fore_color.rgb=WARM;rb.line.color.rgb=TOMATO;rb.line.width=Pt(2)
    tb(s,0.5,1.05,9.0,0.4,"老师说一件事 — 学生用动作回答!",sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.5,1.45,4.4,0.5,"✅ 安全 = 双手举高 (V)",sz=15,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
    tb(s,5.1,1.45,4.4,0.5,"❌ 危险 = 双手交叉 (X)",sz=15,b=True,c=TOMATO,a=PP_ALIGN.CENTER)
    examples=[("🧼","做饭前先洗手","✅"),
              ("🔥","自己开炉火煮东西","❌"),
              ("🔪","把大刀给大人用","✅"),
              ("♨️","自己端一碗热汤","❌"),
              ("🥃","玻璃碎了自己去捡","❌"),
              ("🧽","看到地上有水擦干净","✅")]
    for i,(em,desc,ans) in enumerate(examples):
        col=i%2;row=i//2
        x=0.3+col*4.7;y=2.2+row*0.88
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.5),Inches(0.75))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE
        sh.line.color.rgb=TOMATO if reveal else GRAY;sh.line.width=Pt(1.5)
        tb(s,x+0.12,y+0.15,0.6,0.5,em,sz=22)
        tb(s,x+0.78,y+0.19,3.1,0.4,desc,sz=13,b=True,c=DARK)
        if reveal:
            tb(s,x+3.95,y+0.15,0.5,0.5,ans,sz=22,a=PP_ALIGN.CENTER)
        else:
            tb(s,x+3.95,y+0.15,0.5,0.5,"?",sz=22,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    by=5.02
    bs=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(by),Inches(9.4),Inches(0.42))
    if reveal:
        bs.fill.solid();bs.fill.fore_color.rgb=GREEN_OK;bs.line.fill.background()
        tb(s,0.45,by+0.05,9.0,0.32,"💡 答案揭晓! G1-3 加一句:「___ 不安全，因为 ___」",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    else:
        bs.fill.solid();bs.fill.fore_color.rgb=ORANGE;bs.line.fill.background()
        tb(s,0.45,by+0.05,9.0,0.32,"👉 老师念一句，全班做动作!  ✅ 举手 V  /  ❌ 交叉 X   🤔 3...2...1!",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    notes(s,"玩法(3-5分钟):Question页念6个情境，全班用V/X动作回答；Reveal页揭晓，数对了几个。G1-3加因果句。")
    n+=1;pn(s,n);return s
gesture_game(False)
gesture_game(True)

# ============================================================
# 22 SESSION 2 DIVIDER
# ============================================================
div("Session 2  下午","复习 + 语言目标 (认字 + 写字)\n我会认 7 词 · 我会写 2 词 · 三明治步骤 · 摆餐具",ORANGE,"📖")
n+=1

# ============================================================
# 23 QUICK REVIEW — match danger to safe rule
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔄 快速复习  Quick Review — 危险配安全",ORANGE)
tb(s,0.4,0.85,9,0.3,"把危险和正确做法连起来 (口头)  Match the danger with the safe rule!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
left=[("🔪","刀具",KNIFE),("🔥","炉火",FIRE),("♨️","热水",HOT),("🥃","碎玻璃",GLASS),("💦","湿地面",WET)]
right=["🧍 站着别动叫大人","🙅 不自己开火","🙋 请大人帮忙倒","🧽 擦干净","🤲 交给大人用"]
for i,(em,cn,cl) in enumerate(left):
    y=1.4+i*0.72
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.8),Inches(y),Inches(3.4),Inches(0.58))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.95,y+0.09,3.2,0.4,f"{em} {cn}",sz=16,b=True,c=WHITE)
for i,d in enumerate(right):
    y=1.4+i*0.72
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.8),Inches(y),Inches(3.4),Inches(0.58))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=GREEN_OK;sh.line.width=Pt(2)
    tb(s,5.95,y+0.09,3.2,0.4,d,sz=14,b=True,c=GREEN_OK)
tb(s,4.25,2.95,1.5,0.4,"?",sz=40,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 24-30  我会认 word cards
# ============================================================
def word_card_read(w,py,en,sent,img):
    global n
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",ORANGE)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=FRESH,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=ORANGE,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=ORANGE;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=ORANGE)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    n+=1;pn(s,n);return s
read_words=[
    ("厨房","chú fáng","kitchen","厨房里要小心。","📷 厨房"),
    ("安全","ān quán","safe","用刀不安全，要大人帮忙。","📷 安全"),
    ("洗手","xǐ shǒu","wash hands","做饭前要先洗手。","📷 洗手"),
    ("餐具","cān jù","tableware","我会摆好餐具。","📷 餐具"),
    ("早餐","zǎo cān","breakfast","我给自己准备早餐。","📷 早餐"),
    ("危险","wēi xiǎn","danger","热水很危险。","📷 危险"),
    ("整理","zhěng lǐ","tidy up","吃完饭要整理桌子。","📷 整理"),
]
for w,py,en,sent,img in read_words:
    word_card_read(w,py,en,sent,img)

# ============================================================
# 31 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",ORANGE)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","把字卡贴在\n白板上，老师\n说词语，学生拍！",WARM),
    ("2️⃣","举牌游戏\nShow Me","每人 7 张字卡\n老师说词语\n举正确的卡",RGBColor(0xFF,0xF3,0xE0)),
    ("3️⃣","抢椅子\nMusical Chairs","椅子上放字卡\n音乐停，读出词",RGBColor(0xE8,0xF5,0xE9)),
    ("4️⃣","传话筒\nPass the Mic","传球，停下的人\n读字卡并造句",RGBColor(0xE3,0xF2,0xFD)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 32-33  我会写 cards
# ============================================================
def word_card_write(w,py,en,img):
    global n
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",FRESH)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=FRESH;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=FRESH,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=FRESH)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=FRESH)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    n+=1;pn(s,n);return s
write_words=[
    ("安全","ān quán","safe","📷 安全"),
    ("洗手","xǐ shǒu","wash hands","📷 洗手"),
]
for w,py,en,img in write_words:
    word_card_write(w,py,en,img)

# ============================================================
# 34 SENTENCE FRAMES — 我会说
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💬 我会说  Sentence Frames (K · G1–3)",ORANGE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.55),Inches(4.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ORANGE;sh.line.width=Pt(2.5)
pb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.45),Inches(1.05),Inches(1.6),Inches(0.4))
pb.fill.solid();pb.fill.fore_color.rgb=ORANGE;pb.line.fill.background()
tb(s,0.55,1.10,1.5,0.35,"K (TK-K)",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
k_frames=[("做饭前要先 ____ 。","First, wash ___."),
          ("这个很危险。","This is dangerous."),
          ("我可以帮忙 ____ 。","I can help ___."),
          ("我会洗手。","I can wash hands."),
          ("我会摆餐具。","I can set the table.")]
for i,(cn,en) in enumerate(k_frames):
    y=1.55+i*0.62
    tb(s,0.5,y,4.3,0.4,f"·  {cn}",sz=17,b=True,c=FRESH)
    tb(s,0.5,y+0.32,4.3,0.25,en,sz=9,c=GRAY)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.95),Inches(4.65),Inches(4.0))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=FRESH;sh2.line.width=Pt(2.5)
pb2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.20),Inches(1.05),Inches(1.7),Inches(0.4))
pb2.fill.solid();pb2.fill.fore_color.rgb=FRESH;pb2.line.fill.background()
tb(s,5.30,1.10,1.6,0.35,"G1 - G3",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
g_frames=[("做饭前要先 ___ ，然后 ___ 。","First ___, then ___."),
          ("这个很危险，因为 ___ 。","This is dangerous because ___."),
          ("我可以帮忙 ___ 。","I can help ___."),
          ("吃完饭我要 ___ 。","After eating I will ___.")]
for i,(cn,en) in enumerate(g_frames):
    y=1.65+i*0.78
    tb(s,5.25,y,4.4,0.45,f"·  {cn}",sz=15,b=True,c=FRESH)
    tb(s,5.25,y+0.4,4.4,0.25,en,sz=9,c=GRAY)
tb(s,0.4,5.1,9.2,0.3,"💡 把这张 PPT 截屏打印，贴在每张桌子上 — 学生整堂课参考。",sz=11,b=True,c=TOAST,a=PP_ALIGN.CENTER)
notes(s,"打印当桌签。K重点:短句(先洗手、很危险、帮忙)。G1-3:用因为+然后+完整句。")
pn(s,n)

# ============================================================
# 35 三明治步骤排序 — sandwich step sequencing
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🥪 三明治步骤排序  Order the Steps",FRESH)
tb(s,0.4,0.85,9.2,0.35,"做三明治要按顺序 — 你能排对吗？(先想一想，再看正确顺序)",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[("1️⃣","🧼","洗手"),("2️⃣","🍽️","准备餐具"),("3️⃣","🥬","选择食材"),
       ("4️⃣","🥪","制作三明治"),("5️⃣","🧺","收好食材"),("6️⃣","🧽","清理桌面")]
for i,(num,em,txt) in enumerate(steps):
    col=i%3;row=i//3
    x=0.3+col*3.2;y=1.4+row*1.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=FRESH;sh.line.width=Pt(2.5)
    nb=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.12),Inches(y+0.12),Inches(0.55),Inches(0.55))
    nb.fill.solid();nb.fill.fore_color.rgb=FRESH;nb.line.fill.background()
    tb(s,x+0.12,y+0.16,0.55,0.48,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.18,2.0,0.7,em,sz=34)
    tb(s,x+0.1,y+0.92,2.8,0.4,txt,sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
    if i<5:
        tb(s,x+2.78,y+0.42,0.5,0.5,"➡️",sz=20,c=ORANGE) if col<2 else None
notes(s,"先把6步打乱让孩子排序（可用步骤卡），再展示正确顺序。强调:洗手在最前，清理在最后。")
pn(s,n)

# ============================================================
# 36 早餐搭配 — breakfast components
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🍎 认识早餐搭配  A Balanced Breakfast",ORANGE)
tb(s,0.4,0.85,9.2,0.35,"一份好早餐有 4 样 — 你今天吃了哪些？",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
groups=[("🍞","主食","Grains","面包 · 麦片 · 馒头",TOAST),
        ("🍓","水果/蔬菜","Fruit / Veg","苹果 · 香蕉 · 生菜",FRESH),
        ("🥚","蛋白质","Protein","鸡蛋 · 奶酪 · 豆",ORANGE),
        ("🥛","饮品","Drink","牛奶 · 豆浆 · 水",HOT)]
for i,(em,cn,en,ex,cl) in enumerate(groups):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(2.2),Inches(3.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(2.2),Inches(0.55))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.05,1.42,2.1,0.42,cn,sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.1,2.0,0.9,em,sz=48,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.15,2.0,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    ls=ex.split(" · ")
    tf=tb(s,x+0.1,3.55,2.0,0.4,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=13,c=DARK,a=PP_ALIGN.CENTER)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.9),Inches(9.4),Inches(0.5))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2)
tb(s,0.5,4.97,9.0,0.4,"💬 我的早餐有 ____ 和 ____ 。",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 37 摆餐具 — table setting
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🍽️ 摆餐具练习  Set the Table",FRESH)
tb(s,0.4,0.85,9.2,0.35,"正确摆放盘子、杯子、餐巾和餐具 — 照着摆一摆！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
# Placemat diagram (left)
mat=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.35),Inches(5.2),Inches(3.5))
mat.fill.solid();mat.fill.fore_color.rgb=WARM;mat.line.color.rgb=TOAST;mat.line.width=Pt(2)
plate=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(2.3),Inches(2.3),Inches(1.5),Inches(1.5))
plate.fill.solid();plate.fill.fore_color.rgb=WHITE;plate.line.color.rgb=FRESH;plate.line.width=Pt(2.5)
tb(s,2.3,2.85,1.5,0.4,"🍽️ 盘子",sz=13,b=True,c=FRESH,a=PP_ALIGN.CENTER)
tb(s,1.35,2.6,0.9,0.9,"🍴",sz=28,a=PP_ALIGN.CENTER)
tb(s,1.15,3.35,1.2,0.3,"叉子(左)",sz=10,c=DARK,a=PP_ALIGN.CENTER)
tb(s,3.85,2.6,0.9,0.9,"🥄",sz=28,a=PP_ALIGN.CENTER)
tb(s,3.75,3.35,1.2,0.3,"勺/刀(右)",sz=10,c=DARK,a=PP_ALIGN.CENTER)
tb(s,3.95,1.55,0.9,0.9,"🥛",sz=26,a=PP_ALIGN.CENTER)
tb(s,3.75,2.2,1.3,0.3,"杯子(右上)",sz=10,c=DARK,a=PP_ALIGN.CENTER)
tb(s,1.35,1.6,0.9,0.9,"🧻",sz=26,a=PP_ALIGN.CENTER)
tb(s,1.15,2.25,1.2,0.3,"餐巾(左)",sz=10,c=DARK,a=PP_ALIGN.CENTER)
# Checklist (right)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(1.35),Inches(3.8),Inches(3.5))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=FRESH;sh2.line.width=Pt(2.5)
hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(1.35),Inches(3.8),Inches(0.5))
hd.fill.solid();hd.fill.fore_color.rgb=FRESH;hd.line.fill.background()
tb(s,6.05,1.43,3.5,0.4,"✅ 摆餐具检查",sz=15,b=True,c=WHITE)
tf=tb(s,6.05,2.05,3.6,0.4,"☐ 盘子放中间",sz=15,c=DARK)
for it in ["☐ 叉子放左边","☐ 勺子/刀放右边","☐ 杯子放右上","☐ 餐巾放左边"]:
    ap(tf,"",sz=6);ap(tf,it,sz=15,c=DARK)
tb(s,0.4,4.95,9.2,0.4,"💬 盘子放 ____ ，叉子放 ____ 。",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 38 SESSION 3 DIVIDER
# ============================================================
div("Session 3  下午","动手做项目\n🥪 我给自己做早餐 (三明治)  ·  🎨 我的早餐餐垫",TOAST,"🧑‍🍳")
n+=1

# ============================================================
# 39 Booklet
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,'📓 完成"厨房小帮手"练习册  Day 1 Booklet',TOAST)
ib(s,0.4,0.9,9.2,4.3,"📷 练习册截图 / Booklet pages")
pn(s,n)

# ============================================================
# 40 Projects overview
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 动手时间！  Hands-On Time — 2 个项目",TOAST)
projects=[
    ("PROJECT 1","🥪 我给自己做早餐","My Own Breakfast — Sandwich","洗手 → 选食材 → 做三明治\n→ 介绍 → 收拾干净",WARM,FRESH),
    ("PROJECT 2","🎨 我的早餐餐垫","My Breakfast Placemat","画出盘子、杯子、餐巾、\n餐具的位置 + 检查表",RGBColor(0xFF,0xE9,0xC8),ORANGE),
]
for i,(lbl,nm,en,d,bgc,cl) in enumerate(projects):
    x=0.5+i*4.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(4.3),Inches(4.15))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,4.1,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.4,4.1,0.55,nm,sz=22,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.0,4.1,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,x+0.3,2.45,3.7,1.15,"📷 示范")
    ls=d.split('\n')
    tf=tb(s,x+0.2,3.7,3.95,0.45,ls[0],sz=13,c=DARK,a=PP_ALIGN.CENTER)
    for ln in ls[1:]:ap(tf,ln,sz=13,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 41 Project 1 — Sandwich
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🥪 Project 1: 我给自己做早餐  My Sandwich",FRESH)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.92),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=FRESH;sh.line.fill.background()
tb(s,0.4,0.95,4.2,0.35,"🧺 食材  Ingredients",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.4,4.4,2.0,"🍞 吐司  Toast",sz=13,c=DARK)
ap(tf,"🧀 奶酪  Cheese",sz=13,c=DARK)
ap(tf,"🥬 生菜  Lettuce",sz=13,c=DARK)
ap(tf,"🥕 提前切好的蔬菜  Pre-cut veggies",sz=13,c=DARK)
ap(tf,"🍖 熟食  Cooked deli",sz=13,c=DARK)
# Steps (right)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.92),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=ORANGE;sh2.line.fill.background()
tb(s,5.0,0.95,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.4,4.7,2.0,"① 洗手并清洁桌面",sz=13,c=DARK)
ap(tf2,"② 选择吐司、奶酪、生菜、蔬菜、熟食",sz=13,c=DARK)
ap(tf2,"③ 自己完成三明治",sz=13,c=DARK)
ap(tf2,"④ 介绍:「我的三明治里有 ___ 。」",sz=13,c=DARK)
ap(tf2,"⑤ 收食材、分类垃圾、擦干净桌面",sz=13,c=DARK)
# Safety strip (bottom, full width)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.5),Inches(9.4),Inches(0.95))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=TOMATO;sf.line.width=Pt(2)
tb(s,0.5,3.58,2.0,0.4,"⚠️ 安全要求",sz=14,b=True,c=TOMATO)
tb(s,0.5,3.98,9.2,0.4,"提前确认食物过敏；不用花生酱、生肉、炉火和锋利刀具。",sz=13,b=True,c=DARK)
# Sentence frame
fr=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.6),Inches(9.4),Inches(0.55))
fr.fill.solid();fr.fill.fore_color.rgb=FRESH;fr.line.fill.background()
tb(s,0.5,4.68,9.2,0.4,"🗣️ 我的三明治里有 ____ 和 ____ 。",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 42 Project 2 — Placemat
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎨 Project 2: 我的早餐餐垫  My Placemat",ORANGE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.92),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=FRESH;sh.line.fill.background()
tb(s,0.4,0.95,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.4,4.4,1.5,"📄 A3 纸或卡纸  A3 / cardstock",sz=13,c=DARK)
ap(tf,"🖍️ 彩笔  Markers",sz=13,c=DARK)
ap(tf,"✂️ 剪刀、胶水  Scissors, glue",sz=13,c=DARK)
sh_s=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.05),Inches(4.4),Inches(0.4))
sh_s.fill.solid();sh_s.fill.fore_color.rgb=ORANGE;sh_s.line.fill.background()
tb(s,0.4,3.08,4.2,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf_s=tb(s,0.4,3.5,4.4,1.5,"① 画出盘子、杯子、餐巾、餐具的位置",sz=13,c=DARK)
ap(tf_s,"② 涂上颜色，装饰",sz=13,c=DARK)
ap(tf_s,"③ 加上检查表",sz=13,c=DARK)
# Checklist (right)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.92),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=TOAST;sh2.line.fill.background()
tb(s,5.0,0.95,4.6,0.35,"✅ 餐垫上的检查表  Checklist",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.8,"☐ 我洗手了",sz=16,c=DARK)
for it in ["☐ 我准备餐具了","☐ 我吃完后收拾了","☐ 我擦桌子了"]:
    ap(tf2,"",sz=8);ap(tf2,it,sz=16,c=DARK)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(3.7),Inches(4.8),Inches(1.0))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2)
tb(s,5.0,3.78,4.6,0.35,"🗣️ 展示句型  Say:",sz=13,b=True,c=ORANGE)
tb(s,5.0,4.15,4.6,0.5,"· 盘子放中间，杯子放右上。",sz=13,c=DARK)
pn(s,n)

# ============================================================
# 43 分层 — tiered by age
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌱 分层活动  Tiered by Age",FRESH)
tb(s,0.4,0.9,9.2,0.35,"每个人都能参与 — 按年龄有不同的任务！",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
tiers=[("🐣","低龄学生","Younger","贴图 + 摆放食材\n(照着样子做)",RGBColor(0xD1,0x8F,0x0A)),
       ("🐤","中龄学生","Middle","独立完成三明治\n和餐垫",FRESH),
       ("🐔","高龄学生","Older","设计营养搭配\n+ 帮助低龄同学",ORANGE)]
for i,(em,cn,en,d,cl) in enumerate(tiers):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(3.0),Inches(3.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.6,2.8,0.8,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.5,2.8,0.45,cn,sz=19,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.95,2.8,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.15,3.4,2.7,0.5,ls[0],sz=14,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=14,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 44 回家挑战 — home challenge
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏠 回家挑战  Home Challenge",ORANGE)
tb(s,0.4,0.95,9.2,0.4,"回家试一试，做一个真正的厨房小帮手！",sz=15,b=True,c=FRESH,a=PP_ALIGN.CENTER)
challenges=[("🍽️","帮忙摆餐桌","Help set the table",FRESH),
            ("🥪","准备简单早餐","Make a simple breakfast",ORANGE),
            ("🧽","收拾自己的餐具","Clear your own dishes",TOAST)]
for i,(em,cn,en,cl) in enumerate(challenges):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.6),Inches(3.0),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.85,2.8,0.9,em,sz=48,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.9,2.8,0.5,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.45,2.8,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.35),Inches(9.4),Inches(0.85))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=ORANGE;sf.line.width=Pt(2)
tb(s,0.5,4.45,9.0,0.4,"💬 我在家帮忙 ____ 了！",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,4.83,9.0,0.3,"📸 拍张照片，明天和大家分享! Take a photo and share tomorrow!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 45 Day 1 Badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.3,9,0.7,"🎖️ Day 1 小帮手徽章  Helper Badge",sz=24,b=True,c=FRESH,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.05),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=FRESH;sh.line.width=Pt(5)
tf=tb(s,3.6,1.28,2.8,0.4,"DAY 1",sz=18,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
ap(tf,"🍳",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"厨房小帮手",sz=19,b=True,c=FRESH,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=13,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
ap(tf,"🔪🔥♨️🔌🥃💦",sz=14,a=PP_ALIGN.CENTER)
sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.3),Inches(4.20),Inches(7.4),Inches(0.65))
sb.fill.solid();sb.fill.fore_color.rgb=YOLK;sb.line.color.rgb=ORANGE;sb.line.width=Pt(2.5)
tb(s,1.3,4.25,7.4,0.55,"⭐  ⭐  ⭐  ⭐  ⭐  ⭐",sz=32,b=True,c=TOMATO,a=PP_ALIGN.CENTER)
tb(s,1,4.95,8,0.4,"今天学会了厨房安全 + 做早餐! 🎉",sz=16,b=True,c=FRESH,a=PP_ALIGN.CENTER)
tb(s,1,5.30,8,0.3,"6 个危险 · 红黄绿任务 · 三明治 · 餐垫 · 7 个词",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 46 Tomorrow preview
# ============================================================
s=ns();n+=1;bg(s,FRESH)
tb(s,0.5,0.9,9,0.8,"🍳 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.5,"Day 2 — 待定 (To Be Continued)",sz=26,b=True,c=ORANGE,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🧑‍🍳 继续做厨房小帮手！",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Keep being a kitchen helper!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，小帮手！",sz=15,c=WARM,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/huanli/projects/courseppt/Chinese/厨房小帮手kitchen_helper_pbl/day1_kitchen.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
