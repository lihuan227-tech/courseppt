#!/usr/bin/env python3
"""
小小生活家 Little Life Helper — Day 5 · Session 1: 生活技能嘉年华 Life Skills Celebration
Final class of the unit — review through games + a teacher-led soap-making activity.

Session 1 (约 60 分钟, 40+ 混龄学生, 6–8 桌面小组):
  1. Would You Rather? 生活技能二选一      5–8 分钟   (手势作答, 不离座)
  2. Find the Mistakes 生活技能小侦探      10 分钟    (4 张场景图, 点击逐个揭晓)
  3. Life Skills Charades 你演我猜         10 分钟    (每轮 1 人表演, 组长举手)
  4. Soap-Making 肥皂制作                  约 30 分钟 (老师操作, 8 组各有观察任务)

Reviews Day 1 厨房安全 · Day 2 整理收纳 · Day 3 清洁家务 · Day 4 友善礼貌.
Format mirrors Day 2/3/4 decks: cover · 安排 · 目标 · 分部分隔页 · click-reveal
answers · 教师备注 · 徽章 · 教师准备清单 · 可打印 worksheet & 词卡.
Palette: Carnival (grape + sunny + teal + coral).
16:9 · 10 x 5.625 in · KaiTi/Microsoft YaHei.
Run: python3 create_day5_celebration.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Carnival ---
GRAPE = RGBColor(0x7B, 0x4B, 0xC0)   # primary — celebration purple
SUN   = RGBColor(0xF4, 0xB2, 0x2E)   # sunny yellow
TEAL  = RGBColor(0x16, 0x84, 0x8A)   # unit teal (Day 2)
CORAL = RGBColor(0xEF, 0x6B, 0x53)   # unit coral (Day 2)
ROSE  = RGBColor(0xE8, 0x5B, 0x81)   # Day 4 rose
SKY   = RGBColor(0x3E, 0x8E, 0xC4)
GREEN_OK = RGBColor(0x2E, 0x9E, 0x7A)
RED   = RGBColor(0xD8, 0x45, 0x3A)
GOLD  = RGBColor(0xD1, 0x8F, 0x0A)
CREAM = RGBColor(0xFB, 0xF6, 0xFF)   # light lavender background
WARM  = RGBColor(0xF1, 0xE9, 0xFB)   # lavender panel
MINT  = RGBColor(0xE2, 0xF2, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x2C, 0x2C, 0x2C)
GRAY  = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xBB, 0xBB, 0xBB)
IMGBG = RGBColor(0xEC, 0xE8, 0xF2)

# 四个环节各一色
C_P1 = GRAPE    # Would You Rather
C_P2 = TEAL     # Find the Mistakes
C_P3 = CORAL    # Charades
C_P4 = SUN      # Soap making

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None,fn='KaiTi'):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=fn;return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None,fn='KaiTi'):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=fn
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tf=tb(s,l+0.12,t+0.14,w-0.24,h-0.28,lb,sz=12,c=GRAY,a=PP_ALIGN.CENTER);return tf
def hb(s,txt,c=GRAPE,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE,fn='Microsoft YaHei')
def pn(s,n):
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(9.18),Inches(5.28),Inches(0.5),Inches(0.30))
    chip.fill.solid();chip.fill.fore_color.rgb=WHITE;chip.line.color.rgb=LGRAY;chip.line.width=Pt(0.75)
    bx=s.shapes.add_textbox(Inches(9.18),Inches(5.30),Inches(0.5),Inches(0.26));tf=bx.text_frame;tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    r=p.add_run();r.text=str(n);r.font.size=Pt(10);r.font.color.rgb=GRAY;r.font.name='KaiTi'
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def panel(s,l,t,w,h,fill=WHITE,line=GRAPE,lw=2.5):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line;sh.line.width=Pt(lw)
    return sh
def div(title,sub,color,emoji=""):
    global n
    s=ns();n+=1;bg(s,color)
    tb(s,0.5,1.45,9,1.2,f"{emoji} {title}",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    lines=sub.split("\n")
    tf=tb(s,0.4,2.75,9.2,1.6,lines[0],sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    for ln in lines[1:]:ap(tf,ln,sz=18,c=WHITE,a=PP_ALIGN.CENTER)
    pn(s,n)
    return s
def hint(s,txt,y,color=GRAPE,w=9.2,x=0.4,hgt=0.55,em="💡",sz=14):
    """click-reveal takeaway bar."""
    sh=panel(s,x,y,w,hgt,WARM,color,2.5)
    tb(s,x+0.2,y+(hgt-0.42)/2,w-0.4,0.42,f"{em} {txt}",sz=sz,b=True,c=color)

n=0

# ============================================================
# 1 COVER
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.20,8,0.7,"Life Skills Celebration",sz=28,b=True,c=GRAPE,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,1,0.78,8,0.45,"D5 生活技能嘉年华 · 最后一课",sz=22,b=True,c=GRAPE,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.42),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=GRAPE;sh.line.color.rgb=SUN;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.72),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=SUN;sh2.line.width=Pt(2)
tf=tb(s,3.6,1.94,2.8,0.4,"DAY 5",sz=16,b=True,c=TEAL,a=PP_ALIGN.CENTER)
ap(tf,"🎪🧼",sz=42,a=PP_ALIGN.CENTER)
ap(tf,"生活技能嘉年华",sz=17,b=True,c=GRAPE,a=PP_ALIGN.CENTER)
ap(tf,"THINK · ACT · CREATE",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,4.98,8,0.4,"🥪 🧺 👕 🎒 🧽 💛 🧼  看一看、想一想、演一演、做一做！",sz=14,b=True,c=TEAL,a=PP_ALIGN.CENTER)
notes(s,"Session 1 · 约 60 分钟 · 40+ 名混龄学生。开场语：这是我们「小小生活家」的最后一课！今天不学新本领，而是把这一周学过的厨房安全、整理收纳、清洁家务、友善礼貌，用游戏全部复习一遍，最后一起做肥皂。图片建议：三明治、洗衣机、折好的衣服、书包、海绵和盘子、爱心、肥皂。")
pn(s,n)

# ============================================================
# 2 今天的活动 / 安排
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今天我们要做什么？  What Are We Doing Today?")
items=[("1","二选一  Would You Rather?","5–8 分钟","举手指选一个，说一说为什么",GRAPE),
       ("2","找错误  Find the Mistakes","10 分钟","看图找错，小组轻声讨论",TEAL),
       ("3","你演我猜  Charades","10 分钟","一人表演，组长举手抢答",CORAL),
       ("4","制作肥皂  Make Soap","约 30 分钟","老师操作，每组都有观察任务",SUN)]
for i,(num,nm,tm,dc,cl) in enumerate(items):
    y=0.92+i*1.10
    panel(s,0.5,y,9,0.95,cl,None)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.68),Inches(y+0.20),Inches(0.55),Inches(0.55))
    circ.fill.solid();circ.fill.fore_color.rgb=WHITE;circ.line.fill.background()
    tb(s,0.68,y+0.25,0.55,0.45,num,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.13,4.6,0.4,nm,sz=19,b=True,c=WHITE,fn='Microsoft YaHei')
    tb(s,1.47,y+0.54,4.6,0.32,dc,sz=12.5,c=WHITE)
    tb(s,6.6,y+0.28,2.7,0.4,f"⏱ {tm}",sz=15,b=True,c=WHITE,a=PP_ALIGN.RIGHT)
notes(s,"告诉学生：今天是生活技能主题的最后一课。前三个游戏复习学过的知识，最后一起观察并参与制作肥皂。提醒：全程坐在自己的桌面小组，表现好的小组可以得星星。")
pn(s,n)

# ============================================================
# 3 教学目标
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.80,9,0.4,"📦 内容目标  Content:",sz=17,b=True,c=GRAPE)
tf=tb(s,0.7,1.20,9,1.5,"1. 回顾本周学习的生活技能：厨房安全 · 整理收纳 · 清洁家务 · 友善礼貌",sz=13,c=DARK)
ap(tf,"2. 根据不同生活情境做出合理、安全的选择",sz=13,c=DARK)
ap(tf,"3. 找出厨房、整理、清洁和礼貌行为中的错误，并说出正确做法",sz=13,c=DARK)
ap(tf,"4. 观察并参与简单的肥皂制作过程（融化 → 倒模 → 冷却）",sz=13,c=DARK)
tb(s,0.5,2.72,9,0.4,"🗣️ 语言目标  Language:",sz=17,b=True,c=TEAL)
tb(s,0.7,3.10,4.6,0.9,"👀 我会说：我发现…… / 应该…… / 我选择……",sz=12.5,b=True,c=DARK)
tb(s,5.4,3.10,4.3,0.9,"✍️ 复习词语：洗手 · 分类 · 收拾 · 打扫 · 友善",sz=12.5,b=True,c=DARK)
tb(s,0.7,3.62,9,0.4,"💬 句型：我看到肥皂从______变成了______。我预测肥皂会______。",sz=12.5,b=True,c=DARK)
tb(s,0.5,4.20,9,0.4,"🎨 实践目标：用动作复习生活技能词语 + 参与肥皂制作",sz=13,c=CORAL)
tb(s,0.5,4.62,9,0.4,"🤝 社交目标：合作、等待、倾听、遵守安全规则",sz=13,c=GRAPE)
notes(s,"Students will: review important life skills; make safe and responsible choices; identify mistakes in everyday situations; review vocabulary through actions; observe and participate in a simple soap-making activity; practice teamwork, patience, listening, and safety.")
pn(s,n)

# ============================================================
# 4 我们学过什么
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📚 我们学过什么？  What We Learned",TEAL)
cards=[("🍳","厨房安全与做食物","Kitchen Safety","做饭前洗手\n远离炉灶和热锅\n锅柄朝内 · 小心刀具\n制作三明治",CORAL),
       ("🎒","整理物品","Organizing","折衣服 · 整理书包\n整理行李箱 · 抽屉\n把物品分类摆放",TEAL),
       ("🧽","清洁与家务","Cleaning & Chores","洗碗 · 洗衣服\n分深浅色 · 查口袋\n扫地 · 擦桌子\n物归原位",SKY),
       ("💛","友善与礼貌","Kindness & Manners","请 · 谢谢 · 对不起\n不客气 · 主动帮助\n欢迎客人 · 招待客人",ROSE)]
for i,(em,cn,en,body,cl) in enumerate(cards):
    x=0.42+i*2.32
    panel(s,x,0.90,2.18,3.85,WHITE,cl,2.5)
    tb(s,x,1.02,2.18,0.6,em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x,1.68,2.18,0.4,cn,sz=14.5,b=True,c=cl,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x,2.06,2.18,0.3,en,sz=9.5,c=GRAY,a=PP_ALIGN.CENTER)
    lines=body.split("\n")
    tf=tb(s,x+0.14,2.40,1.92,2.6,"· "+lines[0],sz=11.5,c=DARK)
    for ln in lines[1:]:ap(tf,"· "+ln,sz=11.5,c=DARK)
notes(s,"快速回顾（1–2 分钟）：老师指着每一栏，问「这一天我们学了什么？」让学生齐答关键词即可，不要展开讲解，把时间留给游戏。")
pn(s,n)

# ============================================================
# 5 课堂规则
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"📋 今天的挑战规则  Challenge Rules",GRAPE)
rules=[("👂","认真听","Listen carefully."),("🤔","先思考，再回答","Think before answering."),
       ("🤐","不大声喊答案","Do not shout."),("🤝","小组一起合作","Work with your team."),
       ("✋","用工具前等老师指示","Wait for the teacher."),("🔥","做肥皂不碰热材料","Do not touch hot materials.")]
for i,(em,cn,en) in enumerate(rules):
    x=0.42+(i%3)*3.12; y=0.90+(i//3)*1.18
    panel(s,x,y,2.95,1.05,WHITE,GRAPE,2)
    tb(s,x+0.12,y+0.22,0.6,0.6,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.78,y+0.16,2.05,0.4,cn,sz=13.5,b=True,c=DARK)
    tb(s,x+0.80,y+0.56,2.05,0.35,en,sz=9.5,c=GRAY)
panel(s,0.42,3.32,4.5,1.05,WARM,SUN,2.5)
tb(s,0.60,3.40,4.2,0.35,"✋ 三个手势规则  Hand Signals",sz=13,b=True,c=GOLD)
tf=tb(s,0.60,3.72,4.2,0.6,"👍 我同意    ✋ 我有答案    🤫 小组安静讨论",sz=13,b=True,c=DARK)
panel(s,5.08,3.32,4.5,1.05,WARM,TEAL,2.5)
tb(s,5.26,3.40,4.2,0.35,"🪑 分组  Table Teams",sz=13,b=True,c=TEAL)
tb(s,5.26,3.72,4.2,0.6,"全班分成 6–8 个桌面小组，每组取一个队名，\n合作、安静倾听、答对都可以得星星 ⭐",sz=11.5,c=DARK)
hint(s,"今天不比谁最快，比谁最会合作、最会说清楚！",4.55,GRAPE,em="🌟")
notes(s,"班级人数超过 40 人：全程坐在桌面小组内，不让学生跑动。老师准备一张小组计分表或星星贴纸，奖励合作、安静倾听和完整表达，而不是抢答速度。")
pn(s,n)

# ============================================================
# PART 1 — Would You Rather
# ============================================================
div("Part 1 · 生活技能二选一","Would You Rather?  ⏱ 5–8 分钟\n☝️ 选 A 举一根手指    ✌️ 选 B 举两根手指",C_P1,"🙋")

# 6 游戏说明
s=ns();n+=1;bg(s,CREAM);hb(s,"🙋 怎么玩？  How to Play",C_P1)
steps=[("1️⃣","老师读出两个选择","Teacher reads two choices."),
       ("2️⃣","学生用手指作答，坐在座位上不动","☝️ = A    ✌️ = B"),
       ("3️⃣","请 1–2 名学生说出原因","Use a full sentence."),
       ("4️⃣","不需要每题都解释，选 2–3 题即可","Keep it moving!")]
for i,(em,cn,en) in enumerate(steps):
    y=0.90+i*0.72
    panel(s,0.42,y,5.6,0.64,WHITE,C_P1,2)
    tb(s,0.54,y+0.12,0.5,0.42,em,sz=17,a=PP_ALIGN.CENTER)
    tb(s,1.10,y+0.06,4.8,0.34,cn,sz=13.5,b=True,c=DARK)
    tb(s,1.12,y+0.36,4.8,0.28,en,sz=9.5,c=GRAY)
panel(s,6.20,0.90,3.38,2.80,WARM,SUN,2.5)
tb(s,6.36,0.98,3.06,0.35,"💬 句子提示  Sentence Frames",sz=12.5,b=True,c=GOLD)
tf=tb(s,6.36,1.34,3.06,2.3,"我选择______，因为______。",sz=12,b=True,c=DARK)
ap(tf,"I would choose ___ because ___.",sz=9.5,c=GRAY)
ap(tf,"我觉得______比较容易。",sz=12,b=True,c=DARK)
ap(tf,"I think ___ is easier.",sz=9.5,c=GRAY)
ap(tf,"我比较会______。",sz=12,b=True,c=DARK)
ap(tf,"I am better at ___.",sz=9.5,c=GRAY)
hint(s,"全班坐好，只用手势 —— 不要让学生跑到教室两边！",3.85,C_P1,em="⚠️")
notes(s,"40+ 人大班：坚持手势作答，避免走动。每题 30–45 秒，共 6 题，控制在 5–8 分钟。选择 2–3 题请学生用完整句子说明原因，其他题直接看手势统计人数即可。")
pn(s,n)

# 6 Would-You-Rather 题目
WYR=[("🧺","👕","你愿意折十件衣服，还是洗十个盘子？","Would you rather fold ten shirts or wash ten dishes?","折衣服  Fold clothes","洗盘子  Wash dishes"),
     ("🎒","🗄️","你愿意整理一个乱乱的书包，还是整理一个乱乱的抽屉？","Would you rather organize a messy backpack or a messy drawer?","整理书包  Organize a backpack","整理抽屉  Organize a drawer"),
     ("🥪","🍽️","你愿意做三明治，还是帮忙摆餐具？","Would you rather make a sandwich or set the table?","做三明治  Make a sandwich","摆餐具  Set the table"),
     ("🧹","🧽","你愿意扫地，还是擦桌子？","Would you rather sweep the floor or wipe the desks?","扫地  Sweep the floor","擦桌子  Wipe the desks"),
     ("🧳","🎒","你愿意整理行李箱，还是整理书包？","Would you rather pack a suitcase or organize a backpack?","整理行李箱  Pack a suitcase","整理书包  Organize a backpack"),
     ("🙋","👨‍👩‍👧","你愿意一个人完成家务，还是和家人一起完成？","Would you rather do chores alone or with your family?","一个人  Alone","和家人一起  With my family")]
for i,(emA,emB,cn,en,optA,optB) in enumerate(WYR):
    s=ns();n+=1;bg(s,CREAM);hb(s,f"🙋 二选一 · 第 {i+1} 题  Would You Rather?",C_P1)
    tb(s,0.4,0.86,9.2,0.5,cn,sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,0.4,1.34,9.2,0.35,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    for j,(em,opt,cl,tag) in enumerate([(emA,optA,GRAPE,"A  ☝️"),(emB,optB,TEAL,"B  ✌️")]):
        x=0.60+j*4.55
        panel(s,x,1.80,4.25,2.55,WHITE,cl,3)
        chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.15),Inches(1.92),Inches(1.05),Inches(0.42))
        chip.fill.solid();chip.fill.fore_color.rgb=cl;chip.line.fill.background()
        tb(s,x+0.15,1.95,1.05,0.36,tag,sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,x,2.42,4.25,0.9,em,sz=54,a=PP_ALIGN.CENTER)
        tb(s,x+0.15,3.55,3.95,0.6,opt,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.52,9.2,0.4,"☝️ 选 A 举一根手指      ✌️ 选 B 举两根手指      💬 我选择______，因为______。",
       sz=13,b=True,c=GOLD,a=PP_ALIGN.CENTER)
    notes(s,f"第 {i+1} 题。读题 → 学生举手指 → 老师快速数一数（「A 这边有多少？B 这边有多少？」）。第 2、4、6 题各请 1 名学生用完整句子说明原因即可。图片建议：{optA.split('  ')[0]} / {optB.split('  ')[0]} 的儿童卡通插图。")
    pn(s,n)

# ============================================================
# PART 2 — Find the Mistakes
# ============================================================
div("Part 2 · 生活技能小侦探","Find the Mistakes  ⏱ 10 分钟\n🔍 找错误（厨房 · 书包 · 洗衣）+ 🎭 翻转「剧情」（友善礼貌）",C_P2,"🔍")

# 说明
s=ns();n+=1;bg(s,CREAM);hb(s,"🔍 怎么玩？  Life Skills Detectives",C_P2)
steps=[("1️⃣","看图片，安静观察 20–30 秒","Look carefully. Stay quiet."),
       ("2️⃣","小组轻声讨论 🤫","Whisper with your team."),
       ("3️⃣","把「错误的数量」写在小白板上","Write how many mistakes."),
       ("4️⃣","老师请不同小组说出一个错误","One mistake per team."),
       ("5️⃣","找到错误后，还要说出正确做法","Also say what we SHOULD do.")]
for i,(em,cn,en) in enumerate(steps):
    y=0.90+i*0.66
    panel(s,0.42,y,5.6,0.58,WHITE,C_P2,2)
    tb(s,0.54,y+0.09,0.5,0.4,em,sz=16,a=PP_ALIGN.CENTER)
    tb(s,1.10,y+0.03,4.8,0.32,cn,sz=13,b=True,c=DARK)
    tb(s,1.12,y+0.31,4.8,0.26,en,sz=9.5,c=GRAY)
panel(s,6.20,0.90,3.38,2.60,WARM,SUN,2.5)
tb(s,6.36,0.98,3.06,0.35,"💬 句子提示  Sentence Frames",sz=12.5,b=True,c=GOLD)
tf=tb(s,6.36,1.34,3.06,2.1,"我发现______。",sz=12,b=True,c=DARK)
ap(tf,"I found ___.",sz=9.5,c=GRAY)
ap(tf,"这样不安全，因为______。",sz=12,b=True,c=DARK)
ap(tf,"This is unsafe because ___.",sz=9.5,c=GRAY)
ap(tf,"应该______。",sz=12,b=True,c=DARK)
ap(tf,"We should ___.",sz=9.5,c=GRAY)
hint(s,"每组只说一个错误 —— 让更多小组有机会发言，不重复。",4.30,C_P2,em="🤫")
notes(s,"共 3 张找错误图，每张约 2 分钟；之后进入「翻转剧情」。流程固定：安静观察 → 小组轻声讨论 → 老师点组回答 → 点击揭晓答案。可让每组把错误数量写在小白板上，答对数量的小组加 1 星。")
pn(s,n)

def detective(title_cn,title_en,img_desc,mistakes,fixes,note,color=C_P2):
    """找错误页：左边图片占位，右边错误清单（点击逐个出现）+ 正确做法。"""
    global n
    s=ns();n+=1;bg(s,CREAM);hb(s,f"🔍 {title_cn}",color)
    tb(s,0.4,0.78,9.2,0.3,title_en,sz=11.5,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,0.42,1.10,4.10,3.30,img_desc)
    tb(s,0.42,4.48,4.10,0.4,f"🔍 你能找到 {len(mistakes)} 个错误吗？",sz=14,b=True,c=color,a=PP_ALIGN.CENTER)
    panel(s,4.72,1.10,4.86,1.90,WHITE,RED,2.5)
    tb(s,4.88,1.16,4.5,0.32,"❌ 错误在哪里？（点击逐个出现）",sz=12,b=True,c=RED)
    tf=tb(s,4.88,1.48,4.54,1.45,f"1. {mistakes[0]}",sz=11.5,c=DARK)
    for i,m in enumerate(mistakes[1:]):ap(tf,f"{i+2}. {m}",sz=11.5,c=DARK)
    panel(s,4.72,3.10,4.86,1.90,WHITE,GREEN_OK,2.5)
    tb(s,4.88,3.16,4.5,0.32,"✅ 应该怎么做？",sz=12,b=True,c=GREEN_OK)
    tf2=tb(s,4.88,3.48,4.54,1.45,f"· {fixes[0]}",sz=11.5,c=DARK)
    for f in fixes[1:]:ap(tf2,f"· {f}",sz=11.5,c=DARK)
    notes(s,note);pn(s,n)
    return s

detective("找一找：厨房里有什么错误？","What Is Wrong in the Kitchen?",
    "📷 儿童卡通厨房图（含 5 个明显错误）：\n小朋友没洗手就做饭 · 锅柄朝外 ·\n热汤放在桌子边缘 · 小手伸向刀 ·\n地上有水没有擦干",
    ["小朋友做饭前没有洗手","锅柄朝向外面","热汤放在桌子边上","小朋友把手伸向刀","地上有水，但没有擦干"],
    ["做饭前要先洗手","锅柄应该朝内","热汤要放在桌子中间或靠里面","小朋友不能随便碰刀，要请大人帮忙","地上的水要马上擦干，防止滑倒"],
    "复习 Day 1 厨房安全。先让学生找，再点击揭晓。追问：「为什么锅柄朝外很危险？」（会被碰翻、烫伤）「离炉灶要多远？」（至少一米/三英尺）",CORAL)

detective("找一找：书包里有什么问题？","What Is Wrong With the Backpack?",
    "📷 卡通书包剖面图（含 6 个问题）：\n水杯没盖好 · 水杯压着作业本 ·\n作业纸揉皱 · 香蕉直接放在书上 ·\n书包里有垃圾 · 铅笔散落",
    ["水杯没有盖好","水杯和作业本放在一起","作业纸被揉皱了","香蕉直接放在书本上","书包里有垃圾","铅笔散落，没有放进笔袋"],
    ["盖好、拧紧水杯，放在侧边口袋","水杯和书本分开放","把作业放进文件夹","食物放进饭盒或食品袋","每天清理书包里的垃圾","把文具放进笔袋"],
    "复习 Day 2 整理收纳。可以追问：「大课本应该放在哪个口袋？」（大夹层，靠近背部）「五步整理法第一步是什么？」（全部拿出来）",TEAL)

detective("找一找：洗衣服时有什么错误？","What Is Wrong With the Laundry?",
    "📷 卡通洗衣房图（含 5 个错误）：\n深浅衣服混在一起 · 口袋里有纸巾 ·\n洗衣机塞得太满 · 洗涤标签没看 ·\n地上掉了一只袜子",
    ["白色衣服和深色衣服混在一起","没有检查口袋","洗衣机装得太满","完全没有看洗涤标签","地上有一只袜子没有捡起来"],
    ["深色和浅色衣服分开洗","洗衣服前先检查口袋","洗衣机不能装得太满","看一看衣服上的洗涤标签","把脏衣服全部放进洗衣篮"],
    "复习 Day 3 清洁家务。追问：「口袋里如果有纸巾会怎么样？」（洗完全是碎纸屑）「为什么洗衣机不能塞太满？」（洗不干净）",SKY)

# --- 友善礼貌复习：翻转「剧情」 Flip the Script ---
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 翻转「剧情」怎么玩？  Flip the Script",ROSE)
steps=[("1️⃣","看情景：发生了什么事？","What happened?"),
       ("2️⃣","读一读「冲突对话」⚡","Read the conflict version."),
       ("3️⃣","想一想：听了心里舒服吗？","How would you feel?"),
       ("4️⃣","把它「翻转」成友善的说法 ✨","Flip it into kind words."),
       ("5️⃣","两人一组，读一读、演一演","Read it in pairs.")]
for i,(em,cn,en) in enumerate(steps):
    y=0.90+i*0.66
    panel(s,0.42,y,5.6,0.58,WHITE,ROSE,2)
    tb(s,0.54,y+0.09,0.5,0.4,em,sz=16,a=PP_ALIGN.CENTER)
    tb(s,1.10,y+0.03,4.8,0.32,cn,sz=13,b=True,c=DARK)
    tb(s,1.12,y+0.31,4.8,0.26,en,sz=9.5,c=GRAY)
panel(s,6.20,0.90,3.38,2.60,WARM,SUN,2.5)
tb(s,6.36,0.98,3.06,0.35,"💬 翻转小秘诀  Kind Words",sz=12.5,b=True,c=GOLD)
tf=tb(s,6.36,1.34,3.06,2.1,"我的______被你______了，",sz=12,b=True,c=DARK)
ap(tf,"我感到很______。",sz=12,b=True,c=DARK)
ap(tf,"我希望你______。",sz=12,b=True,c=DARK)
ap(tf,"对不起，我刚刚太______了。",sz=12,b=True,c=DARK)
ap(tf,"你一定很______吧？",sz=12,b=True,c=DARK)
hint(s,"说感受 + 说希望，不说「你干吗……」「你事真多」。",4.30,ROSE,em="✨")
notes(s,"复习 Day 4 友善礼貌。先全班一起做情景一，再让同桌两人一组读情景二、三。重点：把「指责」翻转成「说感受 + 说希望」。时间紧张时只做情景一和情景二。")
pn(s,n)

def flip_scene(idx,title,bg_cn,bg_en,conflict,kind,note):
    """翻转剧情页：上方情景背景，左「冲突对话」右「友善翻转」。"""
    global n
    s=ns();n+=1;bg(s,CREAM);hb(s,f"🎭 翻转「剧情」· {title}  Flip the Script",ROSE)
    panel(s,0.42,0.82,9.16,0.92,WHITE,SUN,2.5)
    tb(s,0.60,0.90,8.8,0.4,f"💡 情景背景：{bg_cn}",sz=14.5,b=True,c=DARK)
    tb(s,0.62,1.32,8.8,0.34,bg_en,sz=10,c=GRAY)
    for j,(head,head_en,em,lines,cl,fill) in enumerate([
            ("冲突对话","Conflict","⚡",conflict,ROSE,RGBColor(0xFD,0xEE,0xF2)),
            ("友善翻转","Kind Flip","✨",kind,GREEN_OK,RGBColor(0xE7,0xF6,0xEF))]):
        x=0.42+j*4.66
        panel(s,x,1.92,4.50,3.05,fill,cl,2.5)
        tb(s,x+0.18,2.00,4.1,0.4,f"{em} {head}",sz=16,b=True,c=cl,fn='Microsoft YaHei')
        tb(s,x+0.20,2.40,4.1,0.3,head_en,sz=10,c=GRAY)
        tf=tb(s,x+0.20,2.74,4.10,2.05,"",sz=6)
        for k,(who,say) in enumerate(lines):
            if k: ap(tf,"",sz=5)
            p=tf.add_paragraph()
            r=p.add_run();r.text=f"{who}：";r.font.size=Pt(13);r.font.bold=True;r.font.color.rgb=cl;r.font.name='KaiTi'
            r2=p.add_run();r2.text=f"「{say}」";r2.font.size=Pt(13);r2.font.color.rgb=DARK;r2.font.name='KaiTi'
    tb(s,0.42,5.02,9.16,0.4,"👉 先读左边，再读右边 —— 哪一种听了心里更舒服？为什么？",sz=13,b=True,c=GOLD,a=PP_ALIGN.CENTER)
    notes(s,note);pn(s,n)
    return s

flip_scene(1,"情景一 · 踩到脚",
    "站路队时，亮亮着急忙慌地挤，一不小心踩到了明明的脚。",
    "Liangliang rushes in line and steps on Mingming's foot.",
    [("明明","你干吗踩我？"),("亮亮","你事真多，我又没看见。")],
    [("明明","我的脚被你踩到了，我感到很疼，希望你站路队时慢一点。"),("亮亮","我刚刚太着急了，没看见。以后我会慢一点。你一定很疼吧，严重吗？")],
    "全班一起做。老师先用生气的语气读左边，再用友善的语气读右边，让学生听出差别。追问：明明在右边先说了什么？（说感受）再说了什么？（说希望）")

flip_scene(2,"情景二 · 借东西",
    "小美想用同桌的彩笔，伸手就去拿。",
    "Xiaomei wants her deskmate's markers and grabs them.",
    [("小美","把那个给我！"),("同桌","凭什么给你？这是我的！")],
    [("小美","请问，我可以用一下你的彩笔吗？我想涂天空。"),("同桌","可以呀，用完记得还给我。"),("小美","谢谢你！")],
    "同桌两人一组读一读，然后请 1–2 组到前面演。重点词：请问 · 可以……吗 · 谢谢你。")

flip_scene(3,"情景三 · 客人来了",
    "奶奶带着客人来家里，乐乐正在玩游戏，头也没抬。",
    "Grandma brings a guest home while Lele is playing a game.",
    [("客人","（站在门口，没有人打招呼……）"),("乐乐","（继续玩，不说话）")],
    [("乐乐","你好，欢迎来我们家！请坐，我给您倒杯水。"),("客人","谢谢你，真有礼貌！")],
    "复习 Day 4 待客礼貌。可以请 2 名学生现场演一遍「开门—打招呼—请坐—倒水」。时间不够时可以只讲不演。")

# ============================================================
# PART 3 — Charades
# ============================================================
div("Part 3 · 生活技能你演我猜","Life Skills Charades  ⏱ 10 分钟\n🎭 只能用动作，不能说话！",C_P3,"🎭")

# 规则
s=ns();n+=1;bg(s,CREAM);hb(s,"🎭 游戏规则  How to Play",C_P3)
rules=[("1️⃣","每组轮流派一名学生到前面"),("2️⃣","学生看一张动作卡（其他人不看）"),
       ("3️⃣","只能用动作表演，不能说话"),("4️⃣","其他小组安静观察，先小组讨论"),
       ("5️⃣","老师说「请举手」，组长举手回答"),("6️⃣","猜对的小组得一颗星 ⭐")]
for i,(em,cn) in enumerate(rules):
    y=0.90+i*0.56
    panel(s,0.42,y,5.6,0.50,WHITE,C_P3,2)
    tb(s,0.54,y+0.06,0.5,0.38,em,sz=15,a=PP_ALIGN.CENTER)
    tb(s,1.10,y+0.05,4.8,0.36,cn,sz=13,b=True,c=DARK)
panel(s,6.20,0.90,3.38,3.36,WARM,GRAPE,2.5)
tb(s,6.36,0.98,3.06,0.35,"👥 40 人以上大班怎么控场",sz=12.5,b=True,c=GRAPE)
tf=tb(s,6.36,1.36,3.06,2.8,"· 每轮只请 1 名学生表演",sz=11.5,c=DARK)
for t in ["· 不允许全班同时喊答案","· 每组先讨论，再由组长举手","· 每个动作控制在 20–30 秒","· 共进行 8–10 轮","· 低年级抽简单词，高年级抽难词"]:
    ap(tf,t,sz=11.5,c=DARK)
hint(s,"抢答不加分，安静等待和答对才加分！",4.40,C_P3,em="⭐")
notes(s,"老师提前把下面两页的词语做成抽取卡（本 PPT 最后一页可直接打印）。混龄班：简单词给低年级，复杂词给高年级。每轮 20–30 秒，8–10 轮共约 10 分钟。")
pn(s,n)

def word_cards(title,words,color,note):
    global n
    s=ns();n+=1;bg(s,CREAM);hb(s,title,color)
    for i,(em,cn,en) in enumerate(words):
        x=0.42+(i%4)*2.32; y=0.95+(i//4)*2.15
        panel(s,x,y,2.18,1.95,WHITE,color,2)
        tb(s,x,y+0.22,2.18,0.7,em,sz=40,a=PP_ALIGN.CENTER)
        tb(s,x,y+1.15,2.18,0.4,cn,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
        tb(s,x,y+1.56,2.18,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    notes(s,note);pn(s,n)
    return s

word_cards("🎭 动作词语 A · 简单（适合低年级）  Easy Words",
    [("🧼","洗手","Wash hands"),("🧹","扫地","Sweep the floor"),("🧽","擦桌子","Wipe the desk"),
     ("👕","折衣服","Fold clothes"),("🍽️","洗碗","Wash dishes"),("🥪","做三明治","Make a sandwich"),
     ("🎒","整理书包","Organize a backpack"),("📦","把物品放回原位","Put things away")],
    C_P3,"低年级抽这一组。老师可以先带全班一起做一遍动作（热身 30 秒），再开始比赛。")

word_cards("🎭 动作词语 B · 挑战（适合高年级）  Challenge Words",
    [("👖","检查口袋","Check the pockets"),("🧳","整理行李箱","Pack a suitcase"),("🗂️","把衣服分类","Sort the clothes"),
     ("🧺","把衣服放进洗衣机","Put clothes in the washer"),("🙋","欢迎客人","Welcome a guest"),("🤝","帮助朋友","Help a friend"),
     ("🙇","说「对不起」","Say sorry"),("🙏","礼貌地请求帮助","Ask politely for help")],
    C_P3,"高年级抽这一组。提示学生：可以分几个动作演出来（先……再……），但全程不能说话、不能用嘴型。")

# 过渡页
s=ns();n+=1;bg(s,CREAM);hb(s,"🌟 我们会这么多生活技能！  We Know So Many Life Skills!",GRAPE)
panel(s,0.42,0.95,4.55,1.75,WHITE,GRAPE,2.5)
tb(s,0.62,1.06,4.2,0.5,"🤔 哪一个生活技能对你最有用？",sz=15,b=True,c=GRAPE)
tb(s,0.62,1.52,4.2,0.35,"Which life skill is most useful to you?",sz=10,c=GRAY)
tb(s,0.62,1.90,4.2,0.6,"💬 我觉得______最有用，因为______。",sz=12.5,b=True,c=DARK)
panel(s,5.13,0.95,4.45,1.75,WHITE,TEAL,2.5)
tb(s,5.33,1.06,4.1,0.5,"🏠 哪一个技能回家后可以继续练习？",sz=15,b=True,c=TEAL)
tb(s,5.33,1.52,4.1,0.35,"Which skill can you practice at home?",sz=10,c=GRAY)
tb(s,5.33,1.90,4.1,0.6,"💬 回家以后，我要______。",sz=12.5,b=True,c=DARK)
tb(s,0.4,2.90,9.2,0.5,"🍳  🎒  🧺  🧽  💛",sz=40,a=PP_ALIGN.CENTER)
hint(s,"只请 2–3 名学生回答，然后进入今天的大活动 —— 做肥皂！",3.65,GRAPE,em="⏭️",sz=15,hgt=0.6)
tb(s,0.4,4.45,9.2,0.5,"🧼 下面，请所有人回到座位，眼睛看老师 —— 我们要开始做肥皂了！",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
notes(s,"过渡页，约 1 分钟。控制人数：只请 2–3 名学生分享。说完立刻进入肥皂环节，趁学生注意力还集中时讲安全规则。")
pn(s,n)

# ============================================================
# PART 4 — 肥皂制作
# ============================================================
div("Part 4 · 今天我们来制作肥皂！","Let's Make Soap!  ⏱ 约 30 分钟\n🧼 老师操作 · 学生观察、闻香味、选颜色、轮流搅拌",C_P4,"🧼")

# 活动介绍
s=ns();n+=1;bg(s,CREAM);hb(s,"🧼 今天我们来制作肥皂！  Let's Make Soap!",C_P4)
tb(s,0.42,0.84,9.2,0.4,"今天的肥皂主要由老师制作，小朋友这样参与：",sz=15,b=True,c=DARK)
parts=[("👀","观察制作步骤","Watch each step"),("👃","闻一闻安全的香味","Smell the fragrance"),
       ("🎨","帮助选择颜色","Vote for the color"),("🥄","在老师指导下轮流搅拌","Take turns stirring"),
       ("🫗","看老师倒进模具","Watch the teacher pour"),("🤔","预测冷却后会怎么样","Predict what happens")]
for i,(em,cn,en) in enumerate(parts):
    x=0.42+(i%3)*3.12; y=1.30+(i//3)*1.20
    panel(s,x,y,2.95,1.05,WHITE,C_P4,2)
    tb(s,x+0.12,y+0.24,0.6,0.6,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.78,y+0.18,2.05,0.4,cn,sz=12.5,b=True,c=DARK)
    tb(s,x+0.80,y+0.58,2.05,0.35,en,sz=9,c=GRAY)
panel(s,0.42,3.78,9.16,1.05,WARM,RED,2.5)
tb(s,0.60,3.86,8.8,0.35,"⚠️ 重要  Important",sz=12.5,b=True,c=RED)
tb(s,0.60,4.16,8.8,0.6,"请根据实际购买的肥皂套装调整步骤。本课默认使用儿童活动常见的「融化后倒模」肥皂套装（Melt & Pour）。\n课堂中不要让学生接触烧碱、强碱或任何腐蚀性材料。",sz=11.5,c=DARK)
notes(s,"Melt & Pour 皂基在微波炉或热水浴中融化即可，全程由老师操作。提前确认是否有学生对香料过敏；香味只加少量。若无加热条件，可提前在办公室融化好，带到教室进行「加色 → 搅拌 → 倒模」环节。")
pn(s,n)

# 安全规则
s=ns();n+=1;bg(s,CREAM);hb(s,"⚠️ 制作肥皂安全规则  Soap-Making Safety Rules",RED)
safety=[("🔥","不碰热的容器","Do not touch hot containers."),
        ("🚫","不把材料放进嘴里","Do not taste the materials."),
        ("✋","没有老师允许，不碰桌上的材料","Wait for the teacher's permission."),
        ("🥄","搅拌时要慢慢搅拌","Stir slowly."),
        ("🫗","倒入模具由老师操作","The teacher pours the soap."),
        ("🧼","完成后要洗手","Wash your hands afterward.")]
for i,(em,cn,en) in enumerate(safety):
    y=0.86+i*0.68
    panel(s,0.42,y,9.16,0.60,WHITE,RED,2)
    tb(s,0.56,y+0.10,0.6,0.45,em,sz=21,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.05,4.6,0.34,f"{i+1}. {cn}",sz=14,b=True,c=DARK)
    tb(s,6.10,y+0.10,3.35,0.34,en,sz=10.5,c=GRAY)
tb(s,0.42,4.94,9.16,0.4,"🙌 老师带读，学生一边做动作一边齐读 —— 读完才发材料！",sz=13,b=True,c=RED,a=PP_ALIGN.CENTER)
notes(s,"必须在发放任何材料之前讲完并齐读。明确划出「教师操作区」：加热设备、热皂液只在这个区域，学生保持距离。")
pn(s,n)

# 材料
s=ns();n+=1;bg(s,CREAM);hb(s,"🧰 我们需要什么？  What Do We Need?",C_P4)
mats=[("🧼","肥皂基","Soap base"),("🎨","安全色素","Soap color"),("🌸","肥皂香味","Fragrance"),("🥄","搅拌棒","Stirring stick"),
      ("🥛","耐热杯","Heat-safe cup"),("🧊","肥皂模具","Soap mold"),("🧤","手套","Gloves"),("🍽️","托盘","Tray")]
for i,(em,cn,en) in enumerate(mats):
    x=0.42+(i%4)*2.32; y=0.90+(i//4)*1.30
    panel(s,x,y,2.18,1.16,WHITE,C_P4,2)
    tb(s,x,y+0.08,2.18,0.5,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.58,2.18,0.32,cn,sz=13.5,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x,y+0.88,2.18,0.26,en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
panel(s,0.42,3.62,9.16,1.30,WARM,TEAL,2.5)
tb(s,0.60,3.70,8.8,0.32,"🧑‍🏫 教师准备  Teacher Prep",sz=12.5,b=True,c=TEAL)
tf=tb(s,0.60,4.02,4.3,0.9,"· 提前把材料按顺序摆好",sz=11.5,c=DARK)
ap(tf,"· 桌面铺桌布或托盘",sz=11.5,c=DARK)
ap(tf,"· 热的材料只放在教师操作区",sz=11.5,c=DARK)
tf2=tb(s,5.10,4.02,4.3,0.9,"· 学生与加热设备保持距离",sz=11.5,c=DARK)
ap(tf2,"· 提前确认学生是否对香味敏感",sz=11.5,c=DARK)
ap(tf2,"· 香味不要添加过多",sz=11.5,c=DARK)
notes(s,"图片要清楚，只展示套装里实际有的材料，不要出现套装中没有的东西，以免学生期待落空。")
pn(s,n)

# 步骤 1–4
s=ns();n+=1;bg(s,CREAM);hb(s,"🧼 肥皂是怎么做出来的？（1–4 步）  How Do We Make Soap?",C_P4)
steps=[("1","🔪","准备肥皂基","Prepare the soap base","老师把肥皂基切成小块",CORAL),
       ("2","🔥","加热融化","Melt the soap base","由老师操作加热设备",RED),
       ("3","🎨","加入颜色","Add color","请学生投票选择颜色",GRAPE),
       ("4","🌸","加入香味","Add fragrance","只加入少量香味材料",ROSE)]
for i,(num,em,cn,en,dc,cl) in enumerate(steps):
    x=0.42+(i%2)*4.60; y=0.92+(i//2)*2.05
    panel(s,x,y,4.45,1.90,WHITE,cl,2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.18),Inches(y+0.16),Inches(0.5),Inches(0.5))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,x+0.18,y+0.21,0.5,0.4,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.80,y+0.14,3.4,0.42,cn,sz=17,b=True,c=cl,fn='Microsoft YaHei')
    tb(s,x+0.82,y+0.56,3.4,0.3,en,sz=10,c=GRAY)
    tb(s,x+3.55,y+0.90,0.8,0.7,em,sz=32,a=PP_ALIGN.CENTER)
    tb(s,x+0.22,y+1.10,3.2,0.6,f"👉 {dc}",sz=12.5,b=True,c=DARK)
notes(s,"第 3 步「选颜色」是全班参与点：给 2–3 个颜色，举手投票，少数服从多数。第 4 步香味只滴 2–3 滴，先让前排闻一闻，再传给「香味观察员」组描述。")
pn(s,n)

# 步骤 5–7
s=ns();n+=1;bg(s,CREAM);hb(s,"🧼 肥皂是怎么做出来的？（5–7 步）  How Do We Make Soap?",C_P4)
steps=[("5","🥄","慢慢搅拌","Stir slowly","学生在老师帮助下轮流搅拌",TEAL),
       ("6","🫗","倒入模具","Pour into the mold","由老师完成，学生站在安全距离观察",SKY),
       ("7","⏳","等待冷却","Let it cool","肥皂冷却后会慢慢变硬",GREEN_OK)]
for i,(num,em,cn,en,dc,cl) in enumerate(steps):
    x=0.42+i*3.12
    panel(s,x,0.92,2.95,2.70,WHITE,cl,2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.18),Inches(1.08),Inches(0.5),Inches(0.5))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,x+0.18,1.13,0.5,0.4,num,sz=18,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.80,1.06,2.0,0.42,cn,sz=16,b=True,c=cl,fn='Microsoft YaHei')
    tb(s,x,1.68,2.95,0.7,em,sz=38,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,2.42,2.65,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,2.75,2.65,0.7,dc,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
hint(s,"倒模时全班安静观察，只有老师可以拿热的容器。",3.80,RED,em="⚠️",sz=15,hgt=0.6)
tb(s,0.42,4.55,9.16,0.7,"⏳ 肥皂大约需要 30–60 分钟才会变硬。如果时间不够，老师可以课后脱模，下次上课或放学时发给小朋友。",
   sz=12.5,b=True,c=GOLD,a=PP_ALIGN.CENTER)
notes(s,"搅拌环节：每组派 1–2 名代表，老师扶住杯子，学生慢慢搅拌 5 秒就换人，避免排长队。倒模由老师完成，学生退后一步观察。")
pn(s,n)

# 学生如何参与（8 组任务）
s=ns();n+=1;bg(s,CREAM);hb(s,"👥 每个人都可以参与  Everyone Can Participate",GRAPE)
tb(s,0.42,0.80,9.2,0.3,"班级人数超过 40 人，不让每个学生都到前面 —— 每组领一个观察任务，全组都有事做。",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
teams=[("1","🎨","颜色观察员","Color Team","观察颜色并投票",GRAPE),
       ("2","👃","香味观察员","Fragrance Team","描述香味，不碰液体",ROSE),
       ("3","📋","步骤专家","Steps Team","把制作步骤说出来",TEAL),
       ("4","🦺","安全监督员","Safety Team","看大家有没有遵守规则",RED),
       ("5","🥄","搅拌助手","Stirring Team","选 1–2 名代表搅拌",SUN),
       ("6","🧊","形状设计师","Shape Team","观察模具的形状",SKY),
       ("7","🔬","科学观察员","Science Team","固体→液体→固体",GREEN_OK),
       ("8","🧹","清洁小助手","Cleaning Team","结束后整理安全物品",CORAL)]
for i,(num,em,cn,en,dc,cl) in enumerate(teams):
    x=0.42+(i%4)*2.32; y=1.16+(i//4)*1.92
    panel(s,x,y,2.18,1.78,WHITE,cl,2)
    chip=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.10),Inches(y+0.10),Inches(0.40),Inches(0.40))
    chip.fill.solid();chip.fill.fore_color.rgb=cl;chip.line.fill.background()
    tb(s,x+0.10,y+0.13,0.40,0.34,num,sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.55,y+0.10,1.55,0.45,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.62,2.18,0.34,cn,sz=13,b=True,c=cl,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x,y+0.94,2.18,0.26,en,sz=8.5,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.12,y+1.22,1.94,0.5,dc,sz=10.5,c=DARK,a=PP_ALIGN.CENTER)
notes(s,"每组只有代表参与实际操作，但全组都要观察并完成本组任务。不要让 40 名学生排队轮流搅拌。老师可以在活动结束时请每组用一句话汇报本组的观察。")
pn(s,n)

# 观察与预测
s=ns();n+=1;bg(s,CREAM);hb(s,"🤔 想一想  Think About It",TEAL)
qs=[("❓","现在的肥皂是固体还是液体？","Is the soap a solid or a liquid now?"),
    ("❄️","冷却以后会发生什么？","What will happen after it cools?"),
    ("🧊","肥皂为什么可以变成模具的形状？","Why does the soap take the shape of the mold?")]
for i,(em,cn,en) in enumerate(qs):
    y=0.90+i*0.90
    panel(s,0.42,y,5.7,0.80,WHITE,TEAL,2)
    tb(s,0.56,y+0.18,0.55,0.45,em,sz=20,a=PP_ALIGN.CENTER)
    tb(s,1.22,y+0.10,4.75,0.36,cn,sz=13.5,b=True,c=DARK)
    tb(s,1.24,y+0.46,4.75,0.3,en,sz=9.5,c=GRAY)
panel(s,6.30,0.90,3.28,2.60,WARM,SUN,2.5)
tb(s,6.46,0.98,2.96,0.32,"💬 句子提示  Frames",sz=12,b=True,c=GOLD)
tf=tb(s,6.46,1.32,2.96,2.1,"我看到肥皂从______变成了______。",sz=11.5,b=True,c=DARK)
ap(tf,"I saw the soap change from ___ to ___.",sz=9,c=GRAY)
ap(tf,"我预测肥皂会______。",sz=11.5,b=True,c=DARK)
ap(tf,"I predict the soap will ___.",sz=9,c=GRAY)
ap(tf,"冷却以后，它会变得______。",sz=11.5,b=True,c=DARK)
ap(tf,"After cooling, it will become ___.",sz=9,c=GRAY)
panel(s,0.42,3.66,9.16,1.20,WARM,GREEN_OK,2.5)
tb(s,0.60,3.74,8.8,0.32,"✅ 参考答案（倒模后点击揭晓）",sz=12,b=True,c=GREEN_OK)
tb(s,0.60,4.06,8.8,0.75,"肥皂基加热后从固体变成液体；倒进模具以后，肥皂液会形成模具的形状；冷却后，它又重新变硬，变回固体。",sz=13,b=True,c=DARK)
notes(s,"在倒入模具之前提问，让学生先预测，再验证。「科学观察员」小组负责回答第三个问题。这是本课的科学连接点：固体—液体—固体的可逆变化。")
pn(s,n)

# 等待时的小活动
s=ns();n+=1;bg(s,CREAM);hb(s,"⏳ 肥皂冷却时，我们可以做什么？  While the Soap Cools",C_P4)
panel(s,0.42,0.92,4.55,3.75,WHITE,GRAPE,2.5)
tb(s,0.60,1.02,4.2,0.4,"选项一 · 设计我的肥皂 🎨",sz=15,b=True,c=GRAPE,fn='Microsoft YaHei')
tb(s,0.60,1.44,4.2,0.3,"Option 1 · Design My Soap",sz=10,c=GRAY)
tb(s,0.60,1.76,4.2,0.5,"在纸上画出：形状 · 颜色 · 香味 · 名字",sz=12,b=True,c=DARK)
tf=tb(s,0.60,2.20,4.2,2.3,"我的肥皂是______形状的。",sz=12,c=DARK)
ap(tf,"My soap is shaped like ___.",sz=9,c=GRAY)
ap(tf,"它是______色的。",sz=12,c=DARK)
ap(tf,"It is ___.",sz=9,c=GRAY)
ap(tf,"它闻起来像______。",sz=12,c=DARK)
ap(tf,"It smells like ___.",sz=9,c=GRAY)
ap(tf,"我给它取名叫______。",sz=12,c=DARK)
ap(tf,"I named it ___.",sz=9,c=GRAY)
panel(s,5.13,0.92,4.45,3.75,WHITE,TEAL,2.5)
tb(s,5.31,1.02,4.1,0.4,"选项二 · 步骤排序 🔢",sz=15,b=True,c=TEAL,fn='Microsoft YaHei')
tb(s,5.31,1.44,4.1,0.3,"Option 2 · Put the Steps in Order",sz=10,c=GRAY)
tb(s,5.31,1.76,4.1,0.4,"打乱顺序的图片，请学生说出正确顺序：",sz=12,b=True,c=DARK)
order=[("🔥","融化"),("🎨","加颜色"),("🌸","加香味"),("🥄","搅拌"),("🫗","倒进模具"),("⏳","等待冷却")]
for i,(em,cn) in enumerate(order):
    x=5.31+(i%3)*1.36; y=2.20+(i//3)*1.05
    panel(s,x,y,1.24,0.92,WARM,TEAL,1.5)
    tb(s,x,y+0.08,1.24,0.4,em,sz=20,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.52,1.24,0.32,cn,sz=11.5,b=True,c=DARK,a=PP_ALIGN.CENTER)
hint(s,"选一个就好 —— 安静、准备简单，正好等肥皂变硬。",4.80,C_P4,em="🕐",hgt=0.5,sz=13)
notes(s,"「设计我的肥皂」worksheet 见本 PPT 倒数第二页，可直接打印。如果肥皂需要更长时间凝固，让学生完成设计纸，老师课后脱模再发放成品。")
pn(s,n)

# ============================================================
# 收尾
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🌟 今天用到了哪些生活技能？  Life Skills We Used Today",GRAPE)
skills=[("👂","认真听"),("⚠️","遵守安全规则"),("⏳","等待轮流"),("🤝","和小组合作"),
        ("🧽","保持桌面整洁"),("🙋","帮助老师"),("💬","礼貌表达意见"),("📦","收拾材料")]
for i,(em,cn) in enumerate(skills):
    x=0.42+(i%4)*2.32; y=0.92+(i//4)*1.05
    panel(s,x,y,2.18,0.92,WHITE,GRAPE,2)
    tb(s,x+0.12,y+0.20,0.55,0.5,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,x+0.72,y+0.26,1.4,0.4,cn,sz=13,b=True,c=DARK)
panel(s,0.42,3.10,9.16,1.25,WARM,TEAL,2.5)
tb(s,0.60,3.20,8.8,0.4,"🧑‍🏫 教师总结  Teacher Summary",sz=12.5,b=True,c=TEAL)
tb(s,0.60,3.56,8.8,0.75,"制作肥皂不只是一个手工活动，我们还用到了安全、合作、耐心、整理和友善。\nMaking soap is not only a craft. We also use safety, teamwork, patience, organization, and kindness.",
   sz=12.5,b=True,c=DARK)
hint(s,"让学生自己说 —— 老师只补充没有被提到的那几项。",4.50,GRAPE,em="💬",hgt=0.55)
notes(s,"这是本课的价值升华：生活技能不只是「会做家务」，更是安全意识、合作、耐心和礼貌。可以给表现突出的小组加星。")
pn(s,n)

# 反思
s=ns();n+=1;bg(s,CREAM);hb(s,"💭 今天我学会了……  Today I Learned…",TEAL)
frames=[("📚","今天我复习了______。","Today I reviewed ___.",GRAPE),
        ("🎉","我最喜欢的活动是______。","My favorite activity was ___.",CORAL),
        ("🏠","我回家以后会帮助家人______。","At home, I will help my family ___.",TEAL),
        ("⚠️","制作肥皂时，我们一定要______。","When making soap, we must ___.",RED)]
for i,(em,cn,en,cl) in enumerate(frames):
    y=0.92+i*0.95
    panel(s,0.42,y,9.16,0.85,WHITE,cl,2)
    tb(s,0.58,y+0.20,0.6,0.5,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.32,y+0.10,7.9,0.4,cn,sz=15,b=True,c=DARK)
    tb(s,1.34,y+0.50,7.9,0.3,en,sz=10,c=GRAY)
tb(s,0.42,4.80,9.16,0.45,"👥 先和旁边的同学互相说一说，再请 3–4 名学生分享 —— 不需要每个人全班发言。",
   sz=13,b=True,c=GOLD,a=PP_ALIGN.CENTER)
notes(s,"Turn and talk：先两两互说 1 分钟，老师巡场听，再点 3–4 名学生（尽量点不同小组、不同年龄段）分享。")
pn(s,n)

# 结束页 / 徽章
s=ns();n+=1;bg(s,CREAM)
hb(s,"🏅 我是生活技能小达人！  I Am a Life Skills Star!",GRAPE)
badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.75),Inches(1.10),Inches(2.9),Inches(2.9))
badge.fill.solid();badge.fill.fore_color.rgb=GRAPE;badge.line.color.rgb=SUN;badge.line.width=Pt(6)
tf=tb(s,0.85,1.45,2.7,0.4,"DAY 5",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
ap(tf,"🎪",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"生活技能嘉年华",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=11,b=True,c=SUN,a=PP_ALIGN.CENTER)
lines=[("🧍","我会照顾自己。","I can take care of myself.",CORAL),
       ("👨‍👩‍👧","我会帮助家人。","I can help my family.",TEAL),
       ("🧹","我会保持整洁。","I can keep things clean and organized.",SKY),
       ("💛","我会友善、有礼貌。","I can be kind and polite.",ROSE),
       ("⚠️","我会安全地学习新技能。","I can learn new skills safely.",GREEN_OK)]
for i,(em,cn,en,cl) in enumerate(lines):
    y=1.10+i*0.78
    panel(s,4.00,y,5.58,0.68,WHITE,cl,2)
    tb(s,4.14,y+0.14,0.5,0.42,em,sz=18,a=PP_ALIGN.CENTER)
    tb(s,4.72,y+0.05,4.7,0.34,cn,sz=14,b=True,c=cl,fn='Microsoft YaHei')
    tb(s,4.74,y+0.37,4.7,0.28,en,sz=9.5,c=GRAY)
tb(s,0.42,4.25,3.5,0.9,"⭐⭐⭐⭐⭐",sz=22,c=SUN,a=PP_ALIGN.CENTER)
tb(s,0.42,4.70,3.5,0.5,"📣 老师：「生活技能——」\n学生：「我最行！」🎉",sz=12,b=True,c=GRAPE,a=PP_ALIGN.CENTER)
notes(s,"结束语与齐呼。图片建议：一群小朋友一起整理、清洁、帮助别人，并展示做好的肥皂。可在此时发放徽章贴纸或肥皂成品。")
pn(s,n)

# ============================================================
# 教师准备清单
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧰 教师课前准备清单  Teacher Prep List",GRAPE)
panel(s,0.42,0.90,4.55,3.10,WHITE,TEAL,2.5)
tb(s,0.60,0.98,4.2,0.35,"🎲 游戏材料  Game Materials",sz=13,b=True,c=TEAL)
tf=tb(s,0.60,1.32,4.2,2.6,"· Would You Rather 题目（本 PPT）",sz=11.5,c=DARK)
for t in ["· 3 张「找错误」场景图（本 PPT 占位处替换）","· 生活技能动作词卡（最后一页可打印）",
          "· 小组计分表或星星贴纸","· 每组一张纸或小白板 + 白板笔","· 分组名单（6–8 组，混龄搭配）"]:
    ap(tf,t,sz=11.5,c=DARK)
panel(s,5.13,0.90,4.45,3.10,WHITE,C_P4,2.5)
tb(s,5.31,0.98,4.1,0.35,"🧼 肥皂制作材料  Soap Materials",sz=13,b=True,c=GOLD)
tf2=tb(s,5.31,1.32,4.1,2.6,"· 肥皂制作套装（Melt & Pour 皂基）",sz=11.5,c=DARK)
for t in ["· 肥皂专用色素 · 香味材料","· 肥皂模具 · 耐热容器 · 搅拌棒","· 手套 · 托盘或桌布 · 纸巾 · 垃圾袋",
          "· 学生设计纸和蜡笔","· 加热设备（只在教师操作区）"]:
    ap(tf2,t,sz=11.5,c=DARK)
panel(s,0.42,4.10,9.16,1.15,WARM,GRAPE,2.5)
tb(s,0.60,4.16,8.8,0.32,"👥 大班课堂管理 10 条  Managing 40+ Students",sz=12,b=True,c=GRAPE)
tb(s,0.60,4.46,4.3,0.75,"1. 全程按桌面小组坐好  2. 二选一只用手势\n3. 找错误先安静观察再讨论  4. 每组只说一个错误\n5. 你演我猜由组长举手作答",sz=10.5,c=DARK)
tb(s,5.10,4.46,4.4,0.75,"6. 设置清楚的「教师操作区」  7. 热材料只由老师碰\n8. 每组分配观察任务  9. 奖励合作与表达，不奖励抢答\n10. 肥皂凝固慢时先做设计纸，课后发成品",sz=10.5,c=DARK)
notes(s,"课前 15 分钟：摆好肥皂材料、划出教师操作区、把动作词卡剪好、每组发一块小白板和笔。确认学生香味过敏情况。")
pn(s,n)

# ============================================================
# 可打印 worksheet · 设计我的肥皂
# ============================================================
s=ns();n+=1;bg(s,WHITE)
hb(s,"🖨️ 可打印 Worksheet · 设计我的肥皂  Design My Soap",GRAPE)
tb(s,0.42,0.80,9.2,0.32,"姓名 Name: ________________        班级 Class: ________________",sz=12,c=DARK)
panel(s,0.42,1.20,4.30,3.90,WHITE,GRAPE,2.5)
tb(s,0.60,1.28,3.9,0.35,"🎨 画一画我的肥皂  Draw your soap",sz=12.5,b=True,c=GRAPE)
panel(s,0.62,1.66,3.90,3.30,WHITE,LGRAY,1.5)
panel(s,4.92,1.20,4.66,3.90,WHITE,TEAL,2.5)
tb(s,5.10,1.28,4.3,0.35,"✍️ 写一写  Write about it",sz=12.5,b=True,c=TEAL)
items=["我的肥皂是__________形状的。  My soap is shaped like ___.",
       "它是____________色的。  It is ______.",
       "它闻起来像____________。  It smells like ______.",
       "我给它取名叫____________。  I named it ______.",
       "我想把它送给__________。  I want to give it to ___."]
for i,t in enumerate(items):
    y=1.72+i*0.66
    tb(s,5.10,y,4.3,0.55,t,sz=10.5,b=True,c=DARK)
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(5.10),Inches(y+0.46),Inches(4.30),Inches(0.012))
    ln.fill.solid();ln.fill.fore_color.rgb=LGRAY;ln.line.fill.background()
tb(s,0.42,5.14,9.16,0.35,"🧼 小小生活家 · Day 5 生活技能嘉年华  Life Skills Celebration",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
notes(s,"打印页：每人一张，肥皂冷却等待时完成。低年级可以只画图 + 说一说；高年级要求写完整句子。")
pn(s,n)

# ============================================================
# 可打印 · Charades 词卡
# ============================================================
s=ns();n+=1;bg(s,WHITE)
hb(s,"✂️ 可打印 · 生活技能动作词卡  Charades Cards (剪开使用)",CORAL)
CARDS=[("🧼","洗手","Wash hands"),("🥪","做三明治","Make a sandwich"),("👕","折衣服","Fold clothes"),
       ("🎒","整理书包","Organize a backpack"),("🧳","整理行李箱","Pack a suitcase"),("🍽️","洗碗","Wash dishes"),
       ("🧽","擦桌子","Wipe the desk"),("🧹","扫地","Sweep the floor"),("👖","检查口袋","Check the pockets"),
       ("🗂️","把衣服分类","Sort the clothes"),("🧺","放进洗衣机","Put clothes in the washer"),("🙋","欢迎客人","Welcome a guest"),
       ("🤝","帮助朋友","Help a friend"),("🙇","说「对不起」","Say sorry"),("📦","把物品放回原位","Put things away")]
for i,(em,cn,en) in enumerate(CARDS):
    x=0.42+(i%5)*1.86; y=0.90+(i//5)*1.46
    panel(s,x,y,1.74,1.34,WHITE,LGRAY,1.5)
    tb(s,x,y+0.10,1.74,0.45,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.62,1.74,0.34,cn,sz=12.5,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x+0.06,y+0.96,1.62,0.3,en,sz=8,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.42,5.22,9.16,0.32,"💡 建议打印在卡纸上并剪开；低年级抽前 8 张，高年级抽后 7 张。",sz=10.5,c=GRAY,a=PP_ALIGN.CENTER)
notes(s,"打印一份即可，老师抽卡使用。也可以背面写编号，方便快速分成「简单/挑战」两叠。")
pn(s,n)

OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "0 final slides ")
if not os.path.isdir(OUT_DIR):
    OUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(OUT_DIR, "day5_life_skills_celebration.pptx")
prs.save(OUT)
print(f"Created {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
