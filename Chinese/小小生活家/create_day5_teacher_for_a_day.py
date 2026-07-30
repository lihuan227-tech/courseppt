#!/usr/bin/env python3
"""
小小生活家 Little Life Helper — Day 5 · 今天你当家！Be the Teacher for a Day
18 slides · 20–30 分钟 · 40 名 6–10 岁混龄学生 · 全班互动为主

Flow: 家里的责任 → 学校的责任 → 小老师领导力 → 今天的任务
  1  封面 今天你当家
  2  本周我们学了什么（Think–Pair–Share）
  3  我们都是生活小帮手（举手回答）
  4  家务打卡表回顾
  5  给自己鼓鼓掌（Call & Response）
  6  过渡：家里有家务……学校呢？
  7  老师每天都要做哪些工作（先猜后揭晓）
  8  老师也有「学校里的家务」（同桌讨论 30 秒）
  9  大揭晓：今天你们就是老师！
  10 一日小老师任务（看 · 提醒 · 榜样 · 报告）
  11–14 四个小组：卫生间 · 图书角 · 书包文具 · 点心区
  15 小老师怎么说（礼貌用语）
  16 小老师守则 DO / DON'T
  17 下午 3:30 大检查
  18 反思 + 结束语

每页备注含：时间 · 教师台词 · 提问 · 学生可能回答 · 40 人管理贴士 · 过渡语。
Palette: Classroom Bright (royal blue + sunny gold + coral + mint).
16:9 · 10 x 5.625 in.
Run: python3 create_day5_teacher_for_a_day.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Classroom Bright ---
ROYAL = RGBColor(0x2F, 0x6F, 0xD6)   # primary blue
GOLD  = RGBColor(0xF5, 0xB0, 0x1F)   # crown gold
CORAL = RGBColor(0xF2, 0x6C, 0x53)
MINT  = RGBColor(0x1F, 0xA9, 0x8F)
GRAPE = RGBColor(0x8A, 0x5A, 0xC8)
ROSE  = RGBColor(0xE8, 0x5B, 0x81)
GREEN_OK = RGBColor(0x2E, 0x9E, 0x7A)
RED   = RGBColor(0xD8, 0x45, 0x3A)
SKY   = RGBColor(0xE8, 0xF2, 0xFD)   # page background
WARM  = RGBColor(0xFF, 0xF3, 0xDC)   # gold panel
SOFT  = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x24, 0x2A, 0x38)
GRAY  = RGBColor(0x83, 0x8A, 0x99)
LGRAY = RGBColor(0xC3, 0xC9, 0xD4)
IMGBG = RGBColor(0xDF, 0xE9, 0xF7)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None,fn='KaiTi'):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=fn;return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None,fn='KaiTi'):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=fn
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def panel(s,l,t,w,h,fill=WHITE,line=ROYAL,lw=2.5):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line;sh.line.width=Pt(lw)
    return sh
def ib(s,l,t,w,h,lb="🖼 插图"):
    """illustration placeholder — replace with a bright child-friendly picture."""
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.color.rgb=LGRAY;sh.line.width=Pt(1.5)
    tb(s,l+0.14,t+0.16,w-0.28,h-0.32,lb,sz=11.5,c=GRAY,a=PP_ALIGN.CENTER)
def hb(s,cn,en,c=ROYAL,t=0.18):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.92))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.45,t+0.06,9.1,0.48,cn,sz=25,b=True,c=WHITE,fn='Microsoft YaHei')
    tb(s,0.47,t+0.56,9.1,0.3,en,sz=13,c=WHITE)
def pn(s,n):
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(9.18),Inches(5.28),Inches(0.5),Inches(0.30))
    chip.fill.solid();chip.fill.fore_color.rgb=WHITE;chip.line.color.rgb=LGRAY;chip.line.width=Pt(0.75)
    bx=s.shapes.add_textbox(Inches(9.18),Inches(5.30),Inches(0.5),Inches(0.26));tf=bx.text_frame;tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    r=p.add_run();r.text=str(n);r.font.size=Pt(10);r.font.color.rgb=GRAY;r.font.name='KaiTi'
def note(s,mins,script,ask,answers,mgmt,nxt):
    """teacher notes block — same 6 fields on every slide."""
    s.notes_slide.notes_text_frame.text = (
        f"⏱ 时间 Time：{mins}\n\n"
        f"🎤 教师台词 Teacher script：{script}\n\n"
        f"❓ 提问 Ask：{ask}\n\n"
        f"💡 学生可能的回答 Expected responses：{answers}\n\n"
        f"👥 40 人班级管理 Management tips：{mgmt}\n\n"
        f"➡️ 过渡语 Transition：{nxt}"
    )
def banner(s,y,txt,color=GOLD,em="🌟",sz=15,h=0.62,fill=WARM):
    panel(s,0.42,y,9.16,h,fill,color,2.5)
    tb(s,0.62,y+(h-0.44)/2,8.8,0.44,f"{em} {txt}",sz=sz,b=True,c=DARK)

n=0

# ============================================================
# 1 COVER
# ============================================================
s=ns();n+=1;bg(s,SKY)
tb(s,0.5,0.28,9,0.9,"👑 今天你当家！",sz=46,b=True,c=ROYAL,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.5,1.18,9,0.5,"Be the Teacher for a Day",sz=22,b=True,c=GOLD,a=PP_ALIGN.CENTER)
ib(s,2.55,1.74,4.9,2.28,"🖼 大插图：一群戴着「小老师」徽章的小朋友，笑着举手\nBright illustration: happy children wearing teacher badges")
tb(s,0.5,4.08,9,0.45,"今天，我们要变成学校的小老师！",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.5,4.50,9,0.32,"Today we become the little teachers of our school!",sz=11.5,c=GRAY,a=PP_ALIGN.CENTER)
banner(s,4.88,"老师问：「准备好了吗？」  学生答：「准备好了！」",ROYAL,"🙋",14,0.5)
note(s,"1 分钟",
     "「小朋友，看老师头上这顶皇冠 👑 —— 今天，你们就是学校的小老师！」先不要解释太多，只要把好奇心点起来。",
     "「准备好了吗？」（重复两次，第二次更大声）",
     "「准备好了！」全班齐答，声音越来越大。",
     "开场用「老师说—学生答」立刻收住 40 人的注意力。要求：手放桌上、眼睛看老师，才开始。",
     "「在当小老师之前，我们先回头看看这一周我们学会了什么本领。」")
pn(s,n)

# ============================================================
# 2 本周我们学了什么
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"📅 本周我们学了什么？","What Did We Learn This Week?",ROYAL)
skills=[("🥪","做三明治","Sandwich making",CORAL),("🍽","洗碗","Dishwashing",ROYAL),
        ("🧺","洗衣服","Laundry",MINT),("🎒","整理书包","Organize backpack",GOLD),
        ("🧹","保持整洁","Keep things clean",GRAPE),("😊","礼貌待人","Be kind & polite",ROSE)]
for i,(em,cn,en,cl) in enumerate(skills):
    x=0.42+(i%3)*3.12; y=1.28+(i//3)*1.42
    panel(s,x,y,2.95,1.28,WHITE,cl,2.5)
    tb(s,x+0.12,y+0.22,0.85,0.8,em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x+1.02,y+0.24,1.8,0.42,cn,sz=17,b=True,c=cl,fn='Microsoft YaHei')
    tb(s,x+1.04,y+0.70,1.85,0.34,en,sz=9.5,c=GRAY)
banner(s,4.28,"和同桌说一说：你最喜欢哪一个？为什么？  Think–Pair–Share (30 秒)",GOLD,"👯",15,0.62)
tb(s,0.42,5.00,9.16,0.35,"💬 我最喜欢______，因为______。  I like ______ because ______.",
   sz=13,b=True,c=ROYAL,a=PP_ALIGN.CENTER)
note(s,"3 分钟",
     "「这一周我们学了六个本领，我们一起念一遍。」老师指图，全班齐读中文，再读英文。然后：「转过身，和你的同桌说一说，你最喜欢哪一个？为什么？」",
     "你最喜欢哪一个活动？为什么？  Which activity did you like best? Why?",
     "「我最喜欢做三明治，因为很好吃。」「我喜欢洗衣服，因为泡泡很好玩。」「我喜欢整理书包，因为找东西很快。」",
     "同桌讨论只给 30 秒，用计时器 + 拍手三下收回注意力。之后只请 3–4 名学生分享，其余用「举手表示你也喜欢这个」代替发言。",
     "「这些本领在家里都用得上 —— 那你们在家真的帮忙做过吗？」")
pn(s,n)

# ============================================================
# 3 我们都是生活小帮手
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"💪 我们都是生活小帮手","We Can Help at Home",MINT)
chores=[("🍽","洗碗"),("👕","叠衣服"),("🧺","分类衣服"),("🥪","做三明治"),
        ("🎒","整理书包"),("🧽","擦桌子"),("🪴","浇花"),("🐶","喂宠物")]
for i,(em,cn) in enumerate(chores):
    x=0.42+(i%4)*2.32; y=1.30+(i//4)*1.32
    panel(s,x,y,2.18,1.18,WHITE,MINT,2)
    tb(s,x,y+0.10,2.18,0.55,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.70,2.18,0.38,cn,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
banner(s,4.10,"举手回答：你在家帮忙做过哪些家务？  Raise your hand!",MINT,"🙋",15,0.62)
tb(s,0.42,4.82,9.16,0.35,"💬 我在家会______。  At home, I can ______.",sz=13,b=True,c=MINT,a=PP_ALIGN.CENTER)
tb(s,0.42,5.16,9.16,0.3,"🖼 建议插图：小朋友洗碗、叠衣服、浇花、喂宠物的彩色插画",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
note(s,"2 分钟",
     "「这些事情，你在家做过哪一件？做过的请举手。」老师逐个念，学生举手，老师大声数「哇，二十个人洗过碗！」最后请 3–4 人说一句完整的话。",
     "你在家帮忙做过哪些家务？  What chores do you do at home?",
     "「我会洗碗。」「我帮妈妈叠衣服。」「我给小狗喂饭。」「我浇花。」",
     "用「举手 + 老师数数」代替逐个发言，40 人也只要 1 分钟。老师要肯定每一个答案：「不管大小，帮忙就很棒！」",
     "「你们做的这些，都记在我们的家务打卡表上 —— 我们一起看看。」")
pn(s,n)

# ============================================================
# 4 家务打卡表回顾
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"📋 家务打卡表回顾","Chore Chart Reflection",GOLD)
ib(s,0.42,1.28,4.30,3.05,"🖼 展示一张家务打卡表样张（贴纸/星星格子）\nSample chore chart with stickers")
qs=[("🌟","谁坚持完成了很多家务？","Who kept going all week?"),
    ("😄","哪一项最容易？","Which one was the easiest?"),
    ("😅","哪一项最难？","Which one was the hardest?")]
for i,(em,cn,en) in enumerate(qs):
    y=1.28+i*1.05
    panel(s,4.92,y,4.66,0.92,WHITE,GOLD,2.5)
    tb(s,5.08,y+0.18,0.55,0.55,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,5.72,y+0.10,3.7,0.4,cn,sz=15,b=True,c=DARK,fn='Microsoft YaHei')
    tb(s,5.74,y+0.50,3.7,0.3,en,sz=9.5,c=GRAY)
banner(s,4.48,"每一次帮助家人，都是在练习责任心。 Helping your family builds responsibility.",GOLD,"🌟",14,0.62)
tb(s,0.42,5.18,9.16,0.3,"👏 老师表扬「坚持」，不表扬「完美」  Celebrate effort, not perfection",
   sz=11,b=True,c=GRAY,a=PP_ALIGN.CENTER)
note(s,"2 分钟",
     "举起打卡表：「这张表上，最重要的不是谁的星星最多，而是谁一直在坚持。」逐条问三个问题，每题只请 1–2 人回答。",
     "谁坚持完成了很多家务？哪一项最容易？哪一项最难？",
     "「洗碗最容易。」「叠衣服最难，衣服总是歪的。」「我忘了两天，但后面又补上了。」",
     "让学生用手势代替发言：容易 👍、难 👎，老师一眼就能看到 40 人的答案。承认「有几天忘记」也要表扬，保护信心。",
     "「你们这一周真的很努力 —— 我们给自己鼓鼓掌！」")
pn(s,n)

# ============================================================
# 5 给自己鼓鼓掌
# ============================================================
s=ns();n+=1;bg(s,GOLD)
tb(s,0.5,0.55,9,0.9,"🎉 给自己鼓鼓掌！",sz=46,b=True,c=WHITE,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.5,1.48,9,0.45,"Celebrate Yourself!",sz=22,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,1.98,9,0.9,"🎊  👏  🎈  ⭐  🎊",sz=44,a=PP_ALIGN.CENTER)
panel(s,1.55,2.92,6.9,1.05,WHITE,WHITE,0)
tb(s,1.70,3.02,6.6,0.45,"这一周，你们已经越来越像真正的小帮手了！",sz=19,b=True,c=GOLD,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,1.70,3.48,6.6,0.35,"You are becoming real life helpers!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
panel(s,1.55,4.15,3.25,0.95,WHITE,WHITE,0)
tb(s,1.65,4.26,3.05,0.4,"老师：「生活小帮手！」",sz=14,b=True,c=ROYAL,a=PP_ALIGN.CENTER)
tb(s,1.65,4.66,3.05,0.34,"Teacher calls…",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
panel(s,5.20,4.15,3.25,0.95,WHITE,WHITE,0)
tb(s,5.30,4.26,3.05,0.4,"学生：「我最棒！」",sz=14,b=True,c=CORAL,a=PP_ALIGN.CENTER)
tb(s,5.30,4.66,3.05,0.34,"Students answer!",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
note(s,"1 分钟",
     "「站起来！我们给自己鼓三下掌 —— 一、二、三！」然后做口号：老师喊「生活小帮手！」，学生答「我最棒！」，连做两遍，第二遍更响。",
     "（不提问，只做庆祝和口号）",
     "「我最棒！」全班齐答 + 三下掌声。",
     "让 40 人一起站起来鼓掌，是很好的「放电」时机。规则先说清楚：只鼓三下，喊完口号马上坐下。老师举手 = 全班安静。",
     "「家里有家务……那学校呢？学校里的事情，是谁在做？」")
pn(s,n)

# ============================================================
# 6 过渡：学校呢？
# ============================================================
s=ns();n+=1;bg(s,SKY)
tb(s,0.5,0.55,9,0.75,"🏠 家里有家务……",sz=34,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.5,1.32,9,0.8,"那学校呢？",sz=40,b=True,c=ROYAL,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.5,2.14,9,0.4,"Who takes care of our school?",sz=17,b=True,c=GOLD,a=PP_ALIGN.CENTER)
tb(s,4.15,2.52,1.7,1.2,"❓",sz=72,a=PP_ALIGN.CENTER)
spots=[("🪑","教室","Classroom",ROYAL),("🚻","卫生间","Bathroom",MINT),
       ("📚","图书角","Books",GRAPE),("🍎","点心区","Snack area",CORAL)]
for i,(em,cn,en,cl) in enumerate(spots):
    x=0.42+i*2.32
    panel(s,x,3.72,2.18,1.28,WHITE,cl,2.5)
    tb(s,x,3.84,2.18,0.5,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x,4.36,2.18,0.36,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x,4.70,2.18,0.28,en,sz=9.5,c=GRAY,a=PP_ALIGN.CENTER)
note(s,"1 分钟",
     "压低声音制造悬念：「在家，是我们帮爸爸妈妈。那在学校呢？教室、卫生间、图书角、点心区 —— 这些地方，是谁在收拾？」停两秒再往下。",
     "学校里的事情，是谁在做？  Who takes care of our school?",
     "「老师！」「保洁阿姨！」「我们自己也可以！」",
     "这一页只用一个问题，不展开讨论。学生喊答案时，老师用手势 🤫 提醒「先想再说」。",
     "「大部分小朋友说是老师 —— 那老师每天到底要做哪些工作呢？」")
pn(s,n)

# ============================================================
# 7 老师每天做什么
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"🧑‍🏫 老师每天都要做哪些工作？","What Do Teachers Do Every Day?",GRAPE)
jobs=[("📚","整理书本"),("🎒","整理书包区"),("✏️","收好铅笔"),("🚻","检查卫生间"),("🍎","清理点心区"),
      ("🧼","擦桌子"),("🪴","浇花"),("🗑","倒垃圾"),("😊","帮助学生"),("❤️","保证大家安全")]
for i,(em,cn) in enumerate(jobs):
    x=0.42+(i%5)*1.86; y=1.30+(i//5)*1.30
    panel(s,x,y,1.74,1.16,WHITE,GRAPE,2)
    tb(s,x,y+0.10,1.74,0.5,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x,y+0.66,1.74,0.38,cn,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
banner(s,4.02,"先猜一猜，再点击揭晓 —— 老师不只是上课，还每天照顾整个教室。",GRAPE,"👀",14,0.62)
tb(s,0.42,4.74,9.16,0.4,"Teachers don't only teach. They take care of the classroom every single day.",
   sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.42,5.12,9.16,0.32,"🖼 建议插图：老师整理书架、擦桌子、检查卫生间的卡通图",sz=10,c=GRAY,a=PP_ALIGN.CENTER)
note(s,"3 分钟",
     "先不显示答案（或用动画遮住）：「猜一猜，老师每天要做多少件事？」学生猜完再逐个点开图标，每揭晓一个就问「这件事你见过吗？见过的举手。」",
     "老师每天都要做哪些工作？你还看到过老师做什么？",
     "「上课」「改作业」「擦白板」「捡地上的铅笔」「拖地」「帮我们找水杯」「送我们上厕所」",
     "40 人一起喊答案会乱：规定「先举手，被点到才说」。老师每揭晓一项就说一句「这也是家务，只不过在学校」，帮助学生建立联系。",
     "「所以你们发现了吗 —— 老师其实也有很多『学校里的家务』。」")
pn(s,n)

# ============================================================
# 8 老师也有学校家务（讨论）
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"🤔 老师是不是也有很多「学校里的家务」？","Do Teachers Have School Chores Too?",CORAL)
ib(s,0.42,1.28,4.30,3.10,"🖼 搞笑插图：一间乱糟糟的教室 —— 书掉一地、\n书包东倒西歪、桌上有点心屑、水杯倒了\nFunny messy classroom illustration")
panel(s,4.92,1.28,4.66,1.55,WARM,CORAL,2.5)
tb(s,5.10,1.40,4.3,0.45,"👯 同桌讨论 30 秒",sz=17,b=True,c=CORAL,fn='Microsoft YaHei')
tb(s,5.12,1.86,4.3,0.34,"Partner talk · 30 seconds",sz=10,c=GRAY)
tb(s,5.10,2.20,4.3,0.5,"如果老师一天都不整理，学校会变成什么样？",sz=14,b=True,c=DARK)
panel(s,4.92,3.00,4.66,1.38,WHITE,GOLD,2.5)
tb(s,5.10,3.10,4.3,0.4,"💬 句子提示  Sentence frames",sz=12.5,b=True,c=GOLD)
tf=tb(s,5.10,3.48,4.3,0.85,"· 教室会变得______。",sz=13,b=True,c=DARK)
ap(tf,"· 我们会找不到______。",sz=13,b=True,c=DARK)
banner(s,4.52,"学校是我们大家的 —— 每个人都可以出一份力！",CORAL,"❤️",14,0.6)
note(s,"2 分钟",
     "「和同桌说一说：如果老师一天都不整理，学校会变成什么样？」放 30 秒计时器。收回后展示乱教室插图，让学生笑一笑，再引导：「所以，学校也需要有人照顾。」",
     "如果老师一天都不整理，学校会变成什么样？  What if no one tidied our school for one day?",
     "「书全掉地上。」「找不到我的水杯。」「地上都是铅笔，会滑倒。」「点心屑会招虫子。」",
     "同桌讨论固定搭配（左边先说，右边后说），30 秒计时器 + 拍手三下收回。只请 2–3 组分享，避免拖时间。",
     "「那今天，我们就来帮老师一个大忙 —— 因为今天……」")
pn(s,n)

# ============================================================
# 9 大揭晓
# ============================================================
s=ns();n+=1;bg(s,ROYAL)
tb(s,0.5,0.60,9,0.85,"今天……",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.5,1.45,9,1.05,"👑 你们就是老师！",sz=52,b=True,c=GOLD,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.5,2.58,9,0.5,"YOU are the teachers today!",sz=24,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ib(s,2.75,3.15,4.5,1.65,"🖼 大插图：小朋友戴着「小老师」徽章，神气地站成一排\nChildren proudly wearing teacher badges")
tb(s,0.5,4.92,9,0.4,"✨ 建议动画：皇冠从上方掉下来 + 掌声音效",sz=11,c=RGBColor(0xCF,0xE0,0xF8),a=PP_ALIGN.CENTER)
note(s,"1 分钟",
     "先卖关子：「今天……（停顿两秒）……你们就是老师！」点击揭晓皇冠。全班欢呼三秒，老师举手示意安静。",
     "「你们准备好当小老师了吗？」",
     "「准备好了！」「耶！」（欢呼）",
     "欢呼一定要「有始有终」：事先说好「欢呼三秒，老师举手就停」。这一页是全课的情绪高点，之后立刻讲任务，把兴奋转成责任感。",
     "「不过，当小老师有一个秘密任务 —— 不是让你们打扫所有东西。」")
pn(s,n)

# ============================================================
# 10 一日小老师任务
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"🎯 一日小老师任务","Today's Mission",GOLD)
tb(s,0.42,1.22,9.16,0.4,"今天的任务不是「把所有东西都打扫干净」，而是：",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
miss=[("👀","观察","Observe","看一看，哪里需要帮忙",ROYAL),
      ("😊","有礼貌地提醒","Politely remind","「请记得……谢谢你！」",MINT),
      ("👍","做好榜样","Be a role model","自己先做到",GOLD),
      ("📝","报告问题","Report problems","告诉老师，不自己处理",CORAL)]
for i,(em,cn,en,dc,cl) in enumerate(miss):
    x=0.42+i*2.32
    panel(s,x,1.72,2.18,2.55,WHITE,cl,2.5)
    tb(s,x,1.88,2.18,0.65,em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x,2.62,2.18,0.42,cn,sz=16,b=True,c=cl,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x,3.06,2.18,0.3,en,sz=9.5,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.14,3.42,1.9,0.7,dc,sz=11.5,c=DARK,a=PP_ALIGN.CENTER)
banner(s,4.45,"用行动做榜样，不用大声命令。  Lead by example, not by yelling.",GOLD,"⭐",15,0.62)
note(s,"2 分钟",
     "「小老师不是『管人的人』，是『帮忙的人』。」逐条读四个任务，每条让全班做一个动作：👀 手放眼睛旁、😊 双手合十、👍 竖大拇指、📝 假装写字。",
     "小老师可以做什么？不可以做什么？",
     "「可以提醒。」「可以帮忙。」「不可以骂人。」「不可以自己去处理打架，要告诉老师。」",
     "四个动作让 40 人一起「用身体记住规则」，比只念一遍有效得多。特别强调第 4 条：遇到不安全的事，报告老师，不自己解决。",
     "「现在，我要给你们分小组 —— 每一组负责学校的一个地方。」")
pn(s,n)

# ============================================================
# 11–14 四个小组
# ============================================================
def group_slide(idx,em,cn,en,dutiesx,img,color,note_args):
    global n
    s=ns();n+=1;bg(s,SKY)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.18),Inches(9.4),Inches(0.92))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.45,0.24,2.2,0.42,f"第{idx}组",sz=17,b=True,c=WHITE,fn='Microsoft YaHei')
    tb(s,0.47,0.66,2.2,0.3,f"Team {idx}",sz=11,c=WHITE)
    tb(s,2.60,0.24,7.0,0.5,f"{em} {cn}",sz=25,b=True,c=WHITE,fn='Microsoft YaHei')
    tb(s,2.62,0.72,7.0,0.3,en,sz=12,c=WHITE)
    ib(s,0.42,1.28,3.60,3.55,img)
    for i,d in enumerate(dutiesx):
        y=1.28+i*0.72
        panel(s,4.22,y,5.36,0.62,WHITE,color,2)
        tb(s,4.38,y+0.11,5.05,0.42,f"✅ {d}",sz=14,b=True,c=DARK)
    tb(s,0.42,5.05,9.16,0.35,"🙂 只做「友善的提醒」，不当「小警察」  Friendly reminders only — never bossy!",
       sz=12,b=True,c=color,a=PP_ALIGN.CENTER)
    note(s,*note_args);pn(s,n)
    return s

group_slide(1,"🚻","卫生间小管家","Bathroom Managers",
    ["提醒大家冲水","纸巾丢进垃圾桶","把台面上的水擦干","看到不安全的行为，报告老师"],
    "🖼 插图：干净明亮的儿童卫生间 —— 洗手台、\n垃圾桶、擦手纸，旁边有「请冲水」小标志",
    MINT,
    ("1.5 分钟",
     "「第一组是卫生间小管家。你们的工作不是打扫厕所，而是提醒和检查。」逐条读，读到「擦干台面」时示范一次擦的动作。",
     "如果地上湿湿的，会发生什么事？",
     "「会滑倒！」「会摔跤。」「要擦干，或者告诉老师。」",
     "明确边界：不碰马桶、不碰脏纸巾，只做提醒和报告。指定 4–5 人一组，其余同学此刻只需要认真听，不要起立。",
     "「第二组，请看这里 —— 你们负责我们的图书角。」"))

group_slide(2,"📚","图书角小管家","Library Managers",
    ["书要放整齐","看完的书放回原位","地上不能有书","爱护图书，不折不画"],
    "🖼 插图：整齐的儿童书架，书脊朝外排好，\n旁边有小朋友把书放回原位",
    GRAPE,
    ("1.5 分钟",
     "「第二组是图书角小管家。书是大家的朋友，我们要照顾它。」示范：把一本书正确地插回书架。",
     "书应该怎么放回去？书掉在地上怎么办？",
     "「书脊朝外。」「放回原来的位置。」「捡起来，放回书架。」",
     "让这一组的同学在座位上做一个「把书插回去」的动作，全班一起做一次，40 人也能参与，不用起身。",
     "「第三组，你们负责的地方每天都最容易乱 —— 书包区！」"))

group_slide(3,"🎒","书包文具小管家","Backpack & Stationery Managers",
    ["书包站好，排整齐","水杯要立起来，盖子拧紧","零食不掉出来","铅笔、彩笔用完放回原处","地上没有铅笔和橡皮"],
    "🖼 插图：整齐的书包柜/挂钩区，书包立好，\n水杯直立，地上干干净净",
    GOLD,
    ("1.5 分钟",
     "「第三组是书包文具小管家。这一周我们学过『五步整理书包』，现在你们要帮全班保持它！」快速回顾：全部拿出来 → 分类 → 放回去。",
     "水杯为什么要立起来？铅笔掉在地上会怎么样？",
     "「不然会漏水，把书打湿。」「会被踩断。」「会滑倒。」",
     "这一组的检查点最多，可以让他们两人一小队分工（书包 / 文具）。提醒：不翻别人的书包，只提醒本人自己整理。",
     "「最后一组 —— 每天最香、也最容易掉渣的地方：点心区！」"))

group_slide(4,"🍎","点心区小管家","Snack Area Managers",
    ["垃圾丢进垃圾桶","桌面擦干净","地上没有食物碎屑","吃完自己收拾好"],
    "🖼 插图：干净的点心桌 —— 桌面擦亮、\n垃圾桶旁没有垃圾、小朋友在擦桌子",
    CORAL,
    ("1.5 分钟",
     "「第四组是点心区小管家。规则很简单：谁吃的，谁收拾。」示范一次「擦桌子三下」的动作，让全班跟着做。",
     "吃完点心，桌上和地上应该是什么样子？",
     "「干干净净。」「垃圾要丢进桶里。」「碎屑要擦掉，不然会有蚂蚁。」",
     "点心时间人多，提醒这一组「排队时先站到旁边等」，避免堵在垃圾桶前。老师同时公布四个组的名单，或用桌号直接分组，节省时间。",
     "「四个组都有任务了 —— 但是，小老师应该怎么说话呢？」"))

# ============================================================
# 15 小老师怎么说
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"💬 小老师怎么说？","Use Polite Words",ROYAL)
panel(s,0.42,1.28,4.50,3.30,RGBColor(0xFD,0xEC,0xEC),RED,2.5)
tb(s,0.62,1.38,4.1,0.45,"❌ 不要这样说",sz=18,b=True,c=RED,fn='Microsoft YaHei')
tb(s,0.64,1.84,4.1,0.3,"Don't say…",sz=10,c=GRAY)
for i,t in enumerate(["「快点！」","「你怎么这样！」","「不许动！」"]):
    y=2.22+i*0.72
    panel(s,0.70,y,3.95,0.58,WHITE,RED,1.5)
    tb(s,0.86,y+0.10,3.7,0.4,t,sz=16,b=True,c=RED,fn='Microsoft YaHei')
panel(s,5.08,1.28,4.50,3.30,RGBColor(0xE9,0xF7,0xEF),GREEN_OK,2.5)
tb(s,5.28,1.38,4.1,0.45,"✅ 请这样说",sz=18,b=True,c=GREEN_OK,fn='Microsoft YaHei')
tb(s,5.30,1.84,4.1,0.3,"Say this instead…",sz=10,c=GRAY)
for i,t in enumerate(["「请记得……」","「谢谢你的合作！」","「我来帮你。」","「请放回原位。」"]):
    y=2.22+i*0.55
    panel(s,5.36,y,3.95,0.45,WHITE,GREEN_OK,1.5)
    tb(s,5.52,y+0.04,3.7,0.36,t,sz=14,b=True,c=GREEN_OK,fn='Microsoft YaHei')
banner(s,4.72,"全班跟老师读一遍：「请记得……谢谢你的合作！」",ROYAL,"🗣",14,0.55)
note(s,"2 分钟",
     "老师先用凶巴巴的语气念左边，再用友善的语气念右边，让学生听出差别。然后全班跟读右边四句，做两遍。",
     "哪一种说法你更愿意听？为什么？",
     "「右边的。」「因为很有礼貌。」「左边的像在骂人。」",
     "跟读时用「老师说一句，学生说一句」的节奏，40 人整齐又不乱。可以请两名学生上台演示一次「提醒同学捡起铅笔」。",
     "「说话要友善 —— 当小老师还有几条重要的守则。」")
pn(s,n)

# ============================================================
# 16 小老师守则
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"📜 小老师守则","Rules for Classroom Managers",MINT)
panel(s,0.42,1.28,4.50,3.55,RGBColor(0xE9,0xF7,0xEF),GREEN_OK,2.5)
tb(s,0.62,1.38,4.1,0.45,"😊 要这样做  DO",sz=18,b=True,c=GREEN_OK,fn='Microsoft YaHei')
for i,(em,t) in enumerate([("💛","友善待人"),("👏","鼓励别人"),("🤝","帮助同学"),("⭐","做好榜样")]):
    y=1.92+i*0.70
    panel(s,0.70,y,3.95,0.58,WHITE,GREEN_OK,1.5)
    tb(s,0.84,y+0.10,0.5,0.4,em,sz=17,a=PP_ALIGN.CENTER)
    tb(s,1.40,y+0.10,3.1,0.4,t,sz=15,b=True,c=DARK,fn='Microsoft YaHei')
panel(s,5.08,1.28,4.50,3.55,RGBColor(0xFD,0xEC,0xEC),RED,2.5)
tb(s,5.28,1.38,4.1,0.45,"❌ 不要这样做  DON'T",sz=18,b=True,c=RED,fn='Microsoft YaHei')
for i,(em,t) in enumerate([("🔊","大声吼叫"),("😤","命令别人"),("😆","嘲笑错误"),("😠","和同学吵架")]):
    y=1.92+i*0.70
    panel(s,5.36,y,3.95,0.58,WHITE,RED,1.5)
    tb(s,5.50,y+0.10,0.5,0.4,em,sz=17,a=PP_ALIGN.CENTER)
    tb(s,6.06,y+0.10,3.1,0.4,t,sz=15,b=True,c=DARK,fn='Microsoft YaHei')
banner(s,4.95,"好的小老师，先管好自己，再帮助别人。",MINT,"👑",14,0.52)
note(s,"1 分钟",
     "「我念左边，你们竖大拇指 👍；我念右边，你们摇摇手 🙅。」用手势快速过一遍八条守则，最后一起说：「先管好自己，再帮助别人。」",
     "如果有同学不听你的提醒，你应该怎么办？",
     "「再说一次，用礼貌的话。」「告诉老师。」「不吵架。」",
     "手势代替发言，1 分钟就能让 40 人全部参与。特别强调「不嘲笑」：小老师看到别人做错，只提醒，不笑话。",
     "「做得好，还有奖励哦 —— 看看下午会发生什么！」")
pn(s,n)

# ============================================================
# 17 大检查
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"🏆 下午 3:30 大检查！","Big Inspection at 3:30 PM",GOLD)
tb(s,0.42,1.22,9.16,0.4,"老师会到每一个区域检查，四个组都有机会拿满星！",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
stars=[("✨","干净","Cleanliness"),("📦","整齐","Organization"),("🙋","责任心","Responsibility"),
       ("💬","有礼貌的提醒","Polite reminders"),("🤝","团队合作","Teamwork")]
for i,(em,cn,en) in enumerate(stars):
    x=0.42+i*1.86
    panel(s,x,1.72,1.74,1.65,WHITE,GOLD,2.5)
    tb(s,x,1.84,1.74,0.5,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x,2.38,1.74,0.36,cn,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x+0.06,2.74,1.62,0.5,en,sz=8.5,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x,3.02,1.74,0.3,"⭐",sz=16,a=PP_ALIGN.CENTER)
panel(s,2.20,3.52,5.60,1.20,WARM,GOLD,3)
tb(s,2.35,3.62,5.3,0.5,"🏆 最佳小老师小组",sz=22,b=True,c=GOLD,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,2.35,4.14,5.3,0.4,"Best Classroom Managers",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.42,4.85,9.16,0.4,"🖼 建议插图：大奖杯 + 四个小组的小朋友一起欢呼",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
note(s,"1 分钟",
     "「下午三点半，老师会拿着这张星星表，去四个区域检查。」逐条念五个评分项，每念一条就问「这一条，你们组做得到吗？」",
     "我们怎样才能拿到五颗星？",
     "「保持干净。」「互相帮忙。」「用礼貌的话提醒。」「大家一起做。」",
     "强调「四个组都可以拿满星」，避免变成互相攀比或告状。可以在黑板上画好四组的星星格子，下午当场贴星星。",
     "「最后，请你想一想：今天，你想成为什么样的小老师？」")
pn(s,n)

# ============================================================
# 18 反思 + 结束
# ============================================================
s=ns();n+=1;bg(s,SKY);hb(s,"💭 今天你想成为什么样的小老师？","What Kind of Little Teacher Will You Be?",GRAPE)
frames=[("🎯","今天我要……","Today I will…",ROYAL),
        ("🙌","我会……","I can…",MINT),
        ("🌟","我希望……","I hope…",GOLD)]
for i,(em,cn,en,cl) in enumerate(frames):
    x=0.42+i*3.12
    panel(s,x,1.28,2.95,1.55,WHITE,cl,2.5)
    tb(s,x,1.42,2.95,0.5,em,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x,1.96,2.95,0.45,cn,sz=19,b=True,c=cl,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
    tb(s,x,2.44,2.95,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
banner(s,2.98,"先和同桌说一说（30 秒），再请 3–4 位小老师分享。",GRAPE,"👯",14,0.58)
panel(s,0.42,3.72,9.16,1.28,WARM,ROSE,3)
tb(s,0.62,3.82,8.8,0.5,"❤️ 一个干净、整齐、友善的学校，需要每一个人的努力！",
   sz=19,b=True,c=ROSE,a=PP_ALIGN.CENTER,fn='Microsoft YaHei')
tb(s,0.62,4.36,8.8,0.4,"Together we make our school a better place.",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.42,5.08,9.16,0.4,"👑 老师：「小老师们，出发！」   学生：「我们准备好了！」",
   sz=14,b=True,c=ROYAL,a=PP_ALIGN.CENTER)
note(s,"2 分钟",
     "「用这三个句子中的一个，和同桌说说你今天要做什么样的小老师。」30 秒后请 3–4 人分享，老师把关键词写在白板上。最后全班一起读红色那句话，喊口号结束。",
     "今天你想成为什么样的小老师？你会怎么帮助大家？",
     "「今天我要提醒大家冲水。」「我会把书放回原位。」「我希望我们组拿五颗星。」",
     "结束语要短、要齐：全班一起读一遍红框句子，再喊一次口号，然后按小组顺序离开，避免 40 人同时起身。",
     "（下课）分发小老师徽章，各组前往自己的区域开始今天的任务。")
pn(s,n)

OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "0 final slides ")
if not os.path.isdir(OUT_DIR):
    OUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(OUT_DIR, "day5 teacher for a day.pptx")
prs.save(OUT)
print(f"Created {OUT}  ({n} slides)")
