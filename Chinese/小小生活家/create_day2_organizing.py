#!/usr/bin/env python3
"""
小小生活家 Little Homemaker — Day 2: 整理和收纳达人 Organizing & Storage Master
Rebuilt to the detailed lesson plan:
  Part 1 老师的混乱书包 (hook)
  Part 2 怎样收拾书包 — 五步整理法 + 检查表
  Part 3 认识收纳好帮手 — 4 工具 (猜->演示->概念) + 配对游戏
  Part 4-5 学习整理行李箱 + 省空间叠衣法 (T恤/裤子) + 摆放挑战
  Part 6 下午任务预告
  语言目标 我会认(5)/我会写(3)/我会说 (kept)
  下午 小组收纳设计挑战
Palette: Fresh Organizing (teal + coral).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette: Fresh Organizing ---
TEAL = RGBColor(0x16,0x84,0x8A)
CORAL = RGBColor(0xEF,0x6B,0x53)
NAVY = RGBColor(0x1E,0x4B,0x54)
SUNNY = RGBColor(0xF4,0xC1,0x3D)
CREAM = RGBColor(0xF8,0xF4,0xEC)
WARM = RGBColor(0xFD,0xEF,0xE6)
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x2C,0x2C,0x2C)
GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
IMGBG = RGBColor(0xEC,0xEC,0xE6)
GREEN_OK = RGBColor(0x2E,0x9E,0x7A)
RED = RGBColor(0xD8,0x45,0x3A)
GOLD = RGBColor(0xD1,0x8F,0x0A)
BLUE = RGBColor(0x3E,0x8E,0xC4)
PURPLE = RGBColor(0x8A,0x6F,0xB0)
AMBER = RGBColor(0xE8,0x9A,0x33)
MINT = RGBColor(0xE2,0xF2,0xF1)

C_VAC = TEAL
C_LAZY = CORAL
C_PEG = BLUE
C_MAT = AMBER

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=TEAL,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n):
    chip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(9.18),Inches(5.28),Inches(0.5),Inches(0.30))
    chip.fill.solid();chip.fill.fore_color.rgb=WHITE;chip.line.color.rgb=LGRAY;chip.line.width=Pt(0.75)
    bx=s.shapes.add_textbox(Inches(9.18),Inches(5.30),Inches(0.5),Inches(0.26));tf=bx.text_frame;tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER
    r=p.add_run();r.text=str(n);r.font.size=Pt(10);r.font.color.rgb=GRAY;r.font.name='KaiTi'
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    lines=sub.split("\n")
    tf=tb(s,0.4,2.75,9.2,1.6,lines[0],sz=22,c=WHITE,a=PP_ALIGN.CENTER)
    for ln in lines[1:]:ap(tf,ln,sz=20,c=WHITE,a=PP_ALIGN.CENTER)
    return s
def step_head(s,num,cn,en,color):
    bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.15),Inches(9.4),Inches(0.72))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.45),Inches(0.24),Inches(0.54),Inches(0.54))
    circ.fill.solid();circ.fill.fore_color.rgb=WHITE;circ.line.fill.background()
    tb(s,0.45,0.28,0.54,0.46,num,sz=20,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,1.2,0.19,8.2,0.42,cn,sz=21,b=True,c=WHITE)
    tb(s,1.22,0.60,8.2,0.24,en,sz=10,c=WARM)

n=0

# ============================================================
# 1 COVER
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,1,0.25,8,0.7,"Organizing & Storage Master",sz=30,b=True,c=TEAL,a=PP_ALIGN.CENTER)
tb(s,1,0.82,8,0.45,"整理和收纳达人",sz=20,c=TEAL,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.25),Inches(1.5),Inches(3.5),Inches(3.5))
sh.fill.solid();sh.fill.fore_color.rgb=TEAL;sh.line.color.rgb=CORAL;sh.line.width=Pt(6)
sh2=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.55),Inches(1.8),Inches(2.9),Inches(2.9))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=CORAL;sh2.line.width=Pt(2)
tf=tb(s,3.6,2.05,2.8,0.4,"DAY 2",sz=16,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"📦",sz=48,a=PP_ALIGN.CENTER)
ap(tf,"整理和收纳",sz=18,b=True,c=TEAL,a=PP_ALIGN.CENTER)
ap(tf,"ORGANIZING & STORAGE",sz=9,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,5.05,8,0.4,"🧹 一起把东西送回家！Let's give everything a home!",sz=14,b=True,c=CORAL,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 2 SCHEDULE
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"⏰ 今日安排  Today's Plan")
for i,(nm,tm,dc,cl) in enumerate([
    ("上午 · 整理书包","Part 1–2","混乱书包 → 五步整理法 → 整理检查表",TEAL),
    ("上午 · 工具 + 行李箱","Part 3–5","收纳好帮手 4 件 + 行李箱省空间叠衣法",CORAL),
    ("下午 · 语言 + 挑战","Part 6","我会认/写/说 + 小组收纳设计挑战",NAVY)]):
    y=0.9+i*1.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9),Inches(1.2))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,0.7,y+0.15,5,0.4,nm,sz=20,b=True,c=WHITE)
    tb(s,0.7,y+0.62,3,0.4,tm,sz=14,c=WARM)
    tb(s,4.5,y+0.38,5.1,0.6,dc,sz=13,c=WHITE)
pn(s,n)

# ============================================================
# 3 OBJECTIVES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎯 教学目标  Learning Objectives")
tb(s,0.5,0.85,9,0.5,"📦 内容目标  Content:",sz=19,b=True,c=TEAL)
tf=tb(s,0.7,1.32,9,1.5,"1. 明白整理的意义：更快找到、减少丢失、节省时间",sz=14,c=DARK)
ap(tf,"2. 学会「五步整理书包」：拿出来 → 做决定 → 分类 → 认识区域 → 放回去",sz=14,c=DARK)
ap(tf,"3. 认识 4 种收纳工具，并为真实问题选对工具",sz=14,c=DARK)
ap(tf,"4. 学会整理行李箱：分类 + 省空间叠衣法 + 合理摆放",sz=14,c=DARK)
tb(s,0.5,3.15,9,0.5,"🗣️ 语言目标  Language:",sz=19,b=True,c=CORAL)
tb(s,0.7,3.6,5.0,0.9,"👀 我会认：混乱 分类 放回\n　　　　　收拾 书包",sz=14,b=True,c=DARK)
tb(s,5.7,3.6,4.0,0.9,"✍️ 我会写：分类 书包 收拾",sz=14,b=True,c=DARK)
tb(s,0.5,4.7,9,0.5,"🎨 实践目标：动手整理书包 + 下午小组收纳设计挑战",sz=14,c=NAVY)
pn(s,n)

# ============================================================
# 4 SESSION 1 DIVIDER
# ============================================================
div("Part 1 · 上午","老师的混乱书包 + 五步整理法\n🎒 为什么找不到  🧺 五步整理  ✅ 检查表",TEAL,"🎒")
n+=1

# ============================================================
# 5 情境导入 — 混乱书包 vs 整理书包 (search game)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔎 为什么总是找不到?  Why Can't I Find It?",CORAL)
tb(s,0.4,0.85,9.2,0.32,"两个书包，装一样的东西 — 老师各找一次，全班计时!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
boxes=[("📦","混乱书包  Messy",RED,["作业本、文件夹、铅笔盒","橡皮、水杯、外套","纸巾、小玩具、废纸","全都混在一起"]),
       ("🎒","整理书包  Tidy",GREEN_OK,["同样的东西","已经分好类","放在固定位置","一下就找到"])]
for i,(em,title,cl,lines) in enumerate(boxes):
    x=0.3+i*4.85
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.3),Inches(4.6),Inches(2.35))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(3)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.3),Inches(4.6),Inches(0.55))
    hd.fill.solid();hd.fill.fore_color.rgb=cl;hd.line.fill.background()
    tb(s,x+0.15,1.37,4.3,0.42,f"{em} {title}",sz=17,b=True,c=WHITE)
    tf=tb(s,x+0.3,1.98,4.1,0.4,f"· {lines[0]}",sz=13,c=DARK)
    for l in lines[1:]:ap(tf,f"· {l}",sz=13,c=DARK)
rb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.35))
rb.fill.solid();rb.fill.fore_color.rgb=WARM;rb.line.color.rgb=CORAL;rb.line.width=Pt(2)
tb(s,0.5,3.93,9.0,0.35,"🎮 玩法  How to Play:",sz=14,b=True,c=CORAL)
tb(s,0.5,4.28,9.2,0.35,"第一次：从混乱书包找橡皮/作业本 (老师边找边说) · 第二次：从整理书包找 · 全班计时!",sz=12,b=True,c=DARK)
tb(s,0.5,4.68,9.2,0.45,"😩 老师台词:「我的作业在哪里?」「怎么这么难找?」「这里怎么还有废纸?」「水杯怎么和作业放一起?」",sz=11,c=GRAY)
notes(s,"准备两个书包(混乱+整理好)，装相同物品：作业本/文件夹/铅笔盒/橡皮/水杯/外套/纸巾/小玩具/废纸。老师先从混乱书包找橡皮或作业本并计时，边找边说台词；再从整理书包找同样物品并计时。把两次时间写白板。")
pn(s,n)

# ============================================================
# 6 观察与讨论 + 教师总结
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💡 观察与讨论  Observe & Discuss",CORAL)
tb(s,0.4,0.85,9.2,0.32,"哪一次更快?为什么?",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
qs=["🕐 哪一次找得更快?","🤔 第一个书包为什么难找?","🚫 哪些东西不该放一起?",
    "📍 怎样整理才能很快找到作业?","⭐ 经常用的东西该放在哪里?"]
tf=tb(s,0.5,1.25,4.5,3.3,qs[0],sz=15,b=True,c=DARK)
for q in qs[1:]:ap(tf,"",sz=7);ap(tf,q,sz=15,b=True,c=DARK)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.25),Inches(4.5),Inches(3.55))
panel.fill.solid();panel.fill.fore_color.rgb=WARM;panel.line.color.rgb=TEAL;panel.line.width=Pt(2.5)
head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.25),Inches(4.5),Inches(0.5))
head.fill.solid();head.fill.fore_color.rgb=TEAL;head.line.fill.background()
tb(s,5.35,1.32,4.2,0.4,"✨ 教师总结  Summary",sz=15,b=True,c=WHITE)
tb(s,5.35,1.9,4.2,0.4,"整理书包不是全部塞进去，而是要:",sz=13,b=True,c=TEAL)
tf2=tb(s,5.4,2.4,4.2,2.2,"🗂️ 分好类",sz=15,b=True,c=DARK)
for t in ["📍 放在固定位置","🔎 容易找到","✋ 容易拿出来","🏠 用完容易放回去"]:
    ap(tf2,"",sz=6);ap(tf2,t,sz=15,b=True,c=DARK)
notes(s,"带全班讨论5个问题，引导说出：混乱=没分类没固定位置。教师总结整理书包5个要点：分好类、固定位置、容易找到、容易拿出、容易放回。")
pn(s,n)

# ============================================================
# 7 第二部分 intro — 怎样收拾书包 (flow)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧺 怎样收拾书包?  How to Pack a Backpack",TEAL)
tb(s,0.4,0.85,9.2,0.32,"先看老师做，再自己动手，最后一起点评!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
flow=[("👩‍🏫","老师示范","跟着 5 个步骤，示范整理一个书包","5 分钟",TEAL),
      ("👧👦","小组动手","5–6 人一组，用书本/笔/水杯一起收拾一个书包","5 分钟",CORAL),
      ("🔎","上台点评","老师拿到前面，和大家一起看哪里做得好","5 分钟",NAVY)]
for i,(em,cn,d,tm,cl) in enumerate(flow):
    y=1.3+i*1.15
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(1.0))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.6),Inches(y+0.2),Inches(0.6),Inches(0.6))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,0.6,y+0.24,0.6,0.5,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.45,y+0.14,3.0,0.4,cn,sz=17,b=True,c=cl)
    tb(s,1.45,y+0.56,6.4,0.4,d,sz=12,c=DARK)
    pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(8.35),Inches(y+0.32),Inches(1.1),Inches(0.36))
    pill.fill.solid();pill.fill.fore_color.rgb=WARM;pill.line.fill.background()
    tb(s,8.35,y+0.35,1.1,0.3,tm,sz=11,b=True,c=cl,a=PP_ALIGN.CENTER)
notes(s,"第二部分流程：(1)老师配合后面的步骤幻灯示范5分钟；(2)每5-6人一组，发书本/笔/水杯等，一起收拾一个书包5分钟；(3)老师拿到前面点评5分钟。")
pn(s,n)

# ============================================================
# 8 第一步 全部拿出来
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"1","把书包里的东西全部拿出来","Take everything out",TEAL)
ib(s,0.3,1.05,4.3,3.05,"📷 把书包倒空 / Empty the backpack")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(3.05))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=TEAL;panel.line.width=Pt(2.5)
tb(s,5.05,1.16,4.5,0.35,"❓ 提问  Ask:",sz=14,b=True,c=TEAL)
tf=tb(s,5.05,1.52,4.5,1.2,"· 东西还在包里，能看清有什么吗?",sz=12,c=DARK)
ap(tf,"· 哪些是今天需要的?",sz=12,c=DARK)
ap(tf,"· 哪些是不需要带的?",sz=12,c=DARK)
qb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(2.75),Inches(4.5),Inches(1.2))
qb.fill.solid();qb.fill.fore_color.rgb=WARM;qb.line.fill.background()
tb(s,5.2,2.84,4.2,0.7,"🗣️ 「整理的第一步不是马上塞回去，\n　　而是先全部拿出来。」",sz=13,b=True,c=CORAL)
gest=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.3),Inches(9.4),Inches(0.55))
gest.fill.solid();gest.fill.fore_color.rgb=TEAL;gest.line.fill.background()
tb(s,0.5,4.38,9.0,0.4,"🙌 全班动作：双手向外张开 —「全部拿出来!」",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"第一步(3分钟)：把书包物品全部放桌面。提问3个。教师语言强调先拿出来。全班做动作：双手向外张开。")
pn(s,n)

# ============================================================
# 9 第二步 检查和做决定 (4 categories)
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"2","检查和做决定","Check & Decide",CORAL)
cats=[("✅","今天需要带","课本·作业本·文件夹·铅笔盒·水杯",GREEN_OK),
      ("🏠","应该放在家里","不需要的小玩具·多余的书·非今天的东西",GOLD),
      ("🗑️","应该丢掉/回收","废纸·空包装·用过的纸巾",RED),
      ("🙋","需要交给大人","学校通知·要签字的表格·不知是谁的东西",BLUE)]
for i,(em,cn,d,cl) in enumerate(cats):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.05+row*1.55
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(1.4))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.18),Inches(y+0.28),Inches(0.8),Inches(0.8))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,x+0.18,y+0.37,0.8,0.6,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,x+1.15,y+0.24,3.25,0.4,cn,sz=16,b=True,c=cl)
    tb(s,x+1.15,y+0.68,3.25,0.6,d,sz=11,c=DARK)
tb(s,0.4,4.28,9.2,0.5,"🤔 问自己：我今天需要它吗?它还能用吗?它该放哪?要交给谁?",sz=13,b=True,c=NAVY,a=PP_ALIGN.CENTER)
notes(s,"第二步(5分钟)：逐一检查物品，分四类——今天需要带/放家里/丢掉回收/交给大人。教师强调四个自问：需要吗、能用吗、放哪、交给谁。")
pn(s,n)

# ============================================================
# 10 第三步 把物品分类 (4 groups)
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"3","把物品分类","Sort into Groups",BLUE)
groups=[("📚","书本类","课本·作业本·阅读书·文件夹",TEAL),
        ("✏️","文具类","铅笔·橡皮·彩笔·尺子·剪刀",CORAL),
        ("🥤","生活用品类","水杯·纸巾·小毛巾·外套",GREEN_OK),
        ("📨","带回家的东西","作业·通知·手工作品",PURPLE)]
for i,(em,cn,d,cl) in enumerate(groups):
    x=0.3+i*2.4
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.05),Inches(2.25),Inches(2.75))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    tb(s,x+0.1,1.2,2.05,0.7,em,sz=36,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.0,2.15,0.4,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    ls=d.split("·")
    tf=tb(s,x+0.15,2.5,1.95,1.2,ls[0],sz=11,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=11,c=DARK,a=PP_ALIGN.CENTER)
qb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.0),Inches(9.4),Inches(0.9))
qb.fill.solid();qb.fill.fore_color.rgb=WARM;qb.line.color.rgb=BLUE;qb.line.width=Pt(2)
tb(s,0.5,4.06,9.0,0.35,"❓ 想一想:",sz=13,b=True,c=BLUE)
tb(s,0.5,4.42,9.2,0.42,"铅笔和橡皮放哪?作业纸能和水杯放一起吗?小玩具属于哪类?哪些要单独放?",sz=12,b=True,c=DARK)
notes(s,"第三步(5分钟)：把桌上物品分四组——书本类/文具类/生活用品类/带回家。提问引导：文具归铅笔盒、干湿分开、玩具归属、单独放的物品。")
pn(s,n)

# ============================================================
# 11 第四步 认识书包区域 (zones)
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"4","认识书包的不同区域","Know the Zones",PURPLE)
ib(s,0.3,1.05,3.5,3.8,"📷 书包分区图 / Backpack zones")
zones=[("🎒","大夹层","课本·文件夹·大作业本 (大书靠背、竖着放)",TEAL),
       ("📓","中间夹层","小作业本·阅读书·练习册·薄文件夹",BLUE),
       ("🧻","前面小口袋","纸巾·小卡片·学生证 (常拿的小物)",PURPLE),
       ("🥤","侧边口袋","水杯 (盖拧紧!)·雨伞",GREEN_OK),
       ("✏️","铅笔盒","铅笔·橡皮·尺子·彩笔 (别散在包底)",CORAL)]
for i,(em,cn,d,cl) in enumerate(zones):
    y=1.05+i*0.76
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.95),Inches(y),Inches(5.75),Inches(0.66))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2)
    tb(s,4.1,y+0.12,0.55,0.45,em,sz=22,a=PP_ALIGN.CENTER)
    tb(s,4.72,y+0.04,1.6,0.35,cn,sz=14,b=True,c=cl)
    tb(s,4.72,y+0.36,4.85,0.28,d,sz=10,c=DARK)
notes(s,"第四步(5分钟)：认识书包区域——大夹层(大书靠背竖放)、中间夹层(薄本子)、前小口袋(纸巾卡片)、侧袋(水杯盖拧紧、雨伞)、铅笔盒(文具别散包底)。")
pn(s,n)

# ============================================================
# 12 第五步 按顺序放回 + 整理原则
# ============================================================
s=ns();n+=1;bg(s,CREAM);step_head(s,"5","按使用顺序放回去","Put Back in Order",GREEN_OK)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(3.75))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=GREEN_OK;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.05),Inches(4.6),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=GREEN_OK;lh.line.fill.background()
tb(s,0.45,1.12,4.4,0.4,"📥 放回顺序  Order",sz=14,b=True,c=WHITE)
order=["1️⃣ 先放最大的书和文件夹","2️⃣ 再放较薄的练习册","3️⃣ 铅笔盒放容易拿的位置","4️⃣ 纸巾、小物进小口袋","5️⃣ 水杯放侧边口袋","6️⃣ 最后检查拉链拉好"]
tf=tb(s,0.5,1.68,4.3,3.0,order[0],sz=13,b=True,c=DARK)
for o in order[1:]:ap(tf,"",sz=5);ap(tf,o,sz=13,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(3.75))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=NAVY;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.05),Inches(4.6),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=NAVY;rh.line.fill.background()
tb(s,5.25,1.12,4.4,0.4,"🔑 整理原则  Rules",sz=14,b=True,c=WHITE)
rules=["⚖️ 重的物品靠近背部","📦 大的放大夹层","🧷 小的放进小口袋","✏️ 文具放进铅笔盒","💧 水杯和书本分开","⭐ 常用的放容易拿到处"]
tf2=tb(s,5.3,1.68,4.3,3.0,rules[0],sz=13,b=True,c=DARK)
for r in rules[1:]:ap(tf2,"",sz=5);ap(tf2,r,sz=13,b=True,c=DARK)
notes(s,"第五步(5分钟)：示范放回顺序6步。整理原则：重的靠背、大的放大夹层、小的进小口袋、文具进铅笔盒、水杯和书分开、常用的易拿。提问：常用书别放最下面、水杯别倒放。")
pn(s,n)

# ============================================================
# 13 书包整理检查表 + 全班口令
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"✅ 书包整理检查表  Backpack Checklist",TEAL)
tb(s,0.4,0.82,9.2,0.3,"和老师一起，一项一项检查:",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
checks=["☐ 没有废纸和垃圾","☐ 书本放得整齐","☐ 作业放进文件夹","☐ 文具放进铅笔盒",
        "☐ 水杯盖紧放侧袋","☐ 每件东西有固定位置","☐ 拉链已经拉好","☐ 书包不会太重"]
for i,ck in enumerate(checks):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.2+row*0.6
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.5))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=TEAL;card.line.width=Pt(1.5)
    tb(s,x+0.15,y+0.08,4.3,0.35,ck,sz=13,b=True,c=DARK)
chant=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.75),Inches(9.4),Inches(1.1))
chant.fill.solid();chant.fill.fore_color.rgb=CORAL;chant.line.fill.background()
tb(s,0.5,3.82,9.0,0.35,"📣 全班口令  Class Chant",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,4.2,9.0,0.35,"老师:「用完东西——」 学生:「送它回家!」",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,4.56,9.0,0.3,"老师:「每天放学——」 学生:「检查书包!」",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
notes(s,"检查表(3分钟)：师生一起逐项检查8项。练习全班口令两句。")
pn(s,n)

# ============================================================
# 14 SESSION DIVIDER — tools + suitcase
# ============================================================
div("Part 3–5 · 上午","收纳好帮手 + 整理行李箱\n🧰 猜工具  🧳 行李箱  👕 省空间叠衣法",CORAL,"🧰")
n+=1

# ============================================================
# 15 猜一猜 overview — 4 tools
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧰 猜一猜：这是什么收纳工具?  Guess the Tool!",TEAL)
tb(s,0.4,0.85,9.2,0.32,"先看工具，别急着说名字 — 观察 → 猜用途 → 演示 → 说它解决了什么问题",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
tools_ov=[("🛍️","真空压缩袋",C_VAC),("🔄","旋转收纳盘",C_LAZY),("🧷","洞洞板",C_PEG),("🧺","玩具收纳垫",C_MAT)]
for i,(em,cn,cl) in enumerate(tools_ov):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.4),Inches(2.25),Inches(2.35))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.63),Inches(1.6),Inches(1.0),Inches(1.0))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,x+0.63,1.73,1.0,0.8,em,sz=40,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,2.15,0.6,cn,sz=15,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.3,2.15,0.3,f"第 {i+1} 个",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,4.05,9.2,0.9,"⭐ 重点展示这 4 种「变化明显、好操作」的工具 — 让学生看得见变化!",sz=14,b=True,c=CORAL,a=PP_ALIGN.CENTER)
notes(s,"第三部分(15-20分钟)：猜一猜教学法——先展示工具不说名字，让学生观察猜用途，再现场演示，最后总结解决了什么问题。重点展示4种变化明显的工具。")
pn(s,n)

# ============================================================
# 16-19 tool cards (guess -> reveal)
# ============================================================
def tool_card(em,cn,en,color,show,guesses,reveal,concept,demo=False,fit=None):
    global n
    s=ns();bg(s,CREAM)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.15),Inches(9.4),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=color;sh.line.fill.background()
    tb(s,0.5,0.22,1.0,0.55,em,sz=28,c=WHITE)
    tb(s,1.5,0.20,5.0,0.5,cn,sz=25,b=True,c=WHITE)
    tb(s,1.5,0.62,5.0,0.25,en,sz=11,c=WARM)
    if demo:
        pill=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.4),Inches(0.27),Inches(2.2),Inches(0.45))
        pill.fill.solid();pill.fill.fore_color.rgb=SUNNY;pill.line.color.rgb=WHITE;pill.line.width=Pt(1.5)
        tb(s,7.5,0.32,2.0,0.4,"⭐ 课堂演示 Demo",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ib(s,0.3,1.05,4.3,2.5,f"📷 {cn}")
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(2.5))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=color;panel.line.width=Pt(2.5)
    tb(s,5.05,1.15,4.5,0.3,"🔎 怎么演示  How to show:",sz=13,b=True,c=color)
    tb(s,5.05,1.46,4.55,0.55,show,sz=11,c=DARK)
    tb(s,5.05,2.12,4.5,0.3,"❓ 让学生猜  Ask kids:",sz=13,b=True,c=color)
    tf=tb(s,5.05,2.42,4.55,0.3,f"· {guesses[0]}",sz=11,c=DARK)
    for g in guesses[1:]:ap(tf,f"· {g}",sz=11,c=DARK)
    rv=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.68),Inches(9.4),Inches(0.82))
    rv.fill.solid();rv.fill.fore_color.rgb=WARM;rv.line.color.rgb=color;rv.line.width=Pt(2)
    tb(s,0.5,3.74,1.2,0.35,"💡 揭晓",sz=13,b=True,c=color)
    tb(s,1.6,3.74,8.0,0.72,reveal,sz=12,b=True,c=DARK)
    kc=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(4.6),Inches(9.4),Inches(0.55))
    kc.fill.solid();kc.fill.fore_color.rgb=color;kc.line.fill.background()
    txt=f"🔑 关键概念: {concept}"
    if fit:txt+=f"   ·   适合: {fit}"
    tb(s,0.5,4.67,9.2,0.4,txt,sz=12,b=True,c=WHITE)
    n+=1;pn(s,n);return s

tool_card("🛍️","真空压缩袋","Vacuum Storage Bag",C_VAC,
    "袋子里放进一件蓬松外套，用手动抽气泵或吸尘器抽走空气，袋子迅速变扁。",
    ["袋子为什么要密封?","圆形接口有什么作用?","抽走空气后会怎样?"],
    "抽走空气，让厚衣服和被子占用更少的空间。",
    "东西没有减少，但占用的空间变小了。",demo=True,fit="厚外套 · 被子 · 枕头 · 换季衣服")
tool_card("🔄","旋转收纳盘","Lazy Susan Turntable",C_LAZY,
    "在托盘上放小瓶子、画笔或调味料，轻轻转动托盘。",
    ["它为什么是圆形的?","怎样拿到放在后面的物品?","它可以放在哪里?"],
    "一转，后面的东西转到前面，不用把其他东西全搬开。",
    "让后面的物品转到前面，更容易拿到。",demo=True,fit="厨房柜 · 冰箱 · 浴室 · 美术区")
tool_card("🧷","洞洞板","Pegboard",C_PEG,
    "一块有很多小洞的板，加上几个挂钩，挂上小物品。",
    ["为什么板上有这么多洞?","挂钩可以换位置吗?","什么物品适合挂起来?"],
    "挂钩可以插在不同位置，根据物品大小自由改变。",
    "利用墙面空间，位置可以根据需要改变。",fit="剪刀 · 胶带 · 小工具 · 耳机 · 美术用品")
tool_card("🧺","玩具收纳垫","Toy Storage Mat",C_MAT,
    "把积木放在垫子上，游戏结束后拉起周围的绳子，收纳垫变成收纳袋。",
    ["为什么周围有绳子?","玩完以后怎样快速整理?","拉动绳子会发生什么?"],
    "在垫子上玩，玩完拉起绳子，玩具就装进袋子里。",
    "好的收纳工具，让整理变得更快更容易。",demo=True,fit="积木 · 拼图 · 小汽车")

# ============================================================
# 20 收纳工具配对游戏
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧩 收纳工具配对游戏  Which Tool Fits?",CORAL)
tb(s,0.4,0.85,9.2,0.32,"老师说一个问题，学生选出最合适的收纳工具!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
matches=[("🧊 厚外套和被子占太多空间","🛍️ 真空压缩袋"),
         ("🍶 柜子后面的小瓶子很难拿到","🔄 旋转收纳盘"),
         ("🧱 墙面有空间，但桌面很乱","🧷 洞洞板"),
         ("🧸 积木玩完要一块块捡起来","🧺 玩具收纳垫")]
for i,(prob,tool) in enumerate(matches):
    y=1.25+i*0.62
    pb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(5.1),Inches(0.52))
    pb.fill.solid();pb.fill.fore_color.rgb=WHITE;pb.line.color.rgb=CORAL;pb.line.width=Pt(1.5)
    tb(s,0.55,y+0.09,4.9,0.4,prob,sz=13,b=True,c=DARK)
    ar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.6),Inches(y),Inches(4.1),Inches(0.52))
    ar.fill.solid();ar.fill.fore_color.rgb=TEAL;ar.line.fill.background()
    tb(s,5.75,y+0.09,3.9,0.4,tool,sz=13,b=True,c=WHITE)
sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.0))
sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=TEAL;sf.line.width=Pt(2.5)
tb(s,0.5,3.92,9.0,0.35,"🧠 选工具时想一想  Choosing a tool:",sz=13,b=True,c=TEAL)
tb(s,0.5,4.28,9.2,0.5,"收纳什么?物品多大?常用吗?放哪最方便?孩子能自己拿放吗?是否真的解决了问题?",sz=12,b=True,c=DARK)
notes(s,"配对游戏(5分钟)：老师念问题，学生选工具。答案：厚外套被子→真空袋；柜后小瓶→旋转盘；墙面空桌面乱→洞洞板；积木→玩具垫。教师总结选工具的6个思考。")
pn(s,n)

# ============================================================
# 21 行李箱装不下了 (情境导入)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧳 行李箱装不下了!  The Suitcase Won't Close!",NAVY)
tb(s,0.4,0.85,9.2,0.32,"老师把衣服随便塞进行李箱 — 一下就满了!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
ib(s,0.3,1.25,4.4,3.05,"📷 塞得乱糟糟的行李箱")
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.25),Inches(4.8),Inches(3.05))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=NAVY;panel.line.width=Pt(2.5)
hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.25),Inches(4.8),Inches(0.5))
hd.fill.solid();hd.fill.fore_color.rgb=NAVY;hd.line.fill.background()
tb(s,5.05,1.32,4.6,0.4,"🕵️ 一起想一想  Discuss",sz=14,b=True,c=WHITE)
tf=tb(s,5.05,1.9,4.5,2.3,"· 为什么行李箱很快就满了?",sz=13,c=DARK)
for q in ["· 什么东西占的空间最大?","· 哪些物品可以放在一起?","· 哪些物品需要分开?","· 怎样更整齐、更省空间?"]:
    ap(tf,"",sz=6);ap(tf,q,sz=13,c=DARK)
tb(s,0.4,4.45,9.2,0.5,"🧳 里面有：T恤 · 裤子 · 外套 · 袜子 · 鞋 · 洗漱用品 · 玩具 · 水杯 · 书",sz=12,b=True,c=NAVY,a=PP_ALIGN.CENTER)
notes(s,"第四部分情境导入(5分钟)：展示打开的行李箱和一堆物品，老师故意随便塞。提问5个引出省空间的需要。")
pn(s,n)

# ============================================================
# 22 整理行李箱基本步骤 (4 steps)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧳 整理行李箱 4 步  4 Steps to Pack",NAVY)
tb(s,0.4,0.85,9.2,0.32,"跟着视频学 — 重点是省空间叠衣法!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
steps=[("1️⃣","先列清单","带真正需要的，不是喜欢的都带",TEAL),
       ("2️⃣","把物品分类","上衣·裤子·内衣袜子·鞋子·洗漱·学习·娱乐",CORAL),
       ("3️⃣","决定摆放位置","重的靠轮子·衣服折或卷·鞋子装袋·洗漱防水袋·常用易拿",BLUE),
       ("4️⃣","最后检查","漏带?带太多?液体密封?关得上?会滚动?",GREEN_OK)]
for i,(num,cn,d,cl) in enumerate(steps):
    y=1.25+i*0.9
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.78))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(2.5)
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.6),Inches(y+0.13),Inches(0.52),Inches(0.52))
    circ.fill.solid();circ.fill.fore_color.rgb=cl;circ.line.fill.background()
    tb(s,0.6,y+0.16,0.52,0.44,num,sz=17,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1.35,y+0.09,2.6,0.4,cn,sz=16,b=True,c=cl)
    tb(s,3.95,y+0.22,5.6,0.4,d,sz=12,c=DARK)
notes(s,"整理行李箱基本步骤(5分钟，根据视频)：1先列清单(带需要的)；2分类(上衣/裤子/内衣袜子/鞋/洗漱/学习/娱乐)；3摆放位置(重的靠轮子、衣服折或卷、鞋装袋、洗漱防水袋、小物小袋、常用易拿)；4最后检查(漏带/带太多/液体密封/关得上/会滚动)。")
pn(s,n)

# ============================================================
# 23 省空间叠T恤 (卷衣法 + 比较实验)
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👕 省空间叠 T 恤：卷衣法  Roll a T-shirt",CORAL)
tb(s,0.4,0.85,9.2,0.32,"卷起来更省空间、不容易皱 — 跟着做!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.6),Inches(3.5))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=CORAL;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.6),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=CORAL;lh.line.fill.background()
tb(s,0.45,1.32,4.4,0.4,"🌯 卷衣法五步  5 Steps",sz=14,b=True,c=WHITE)
rsteps=["1️⃣ T恤正面朝下铺平","2️⃣ 两边向中间折","3️⃣ 把袖子整理好","4️⃣ 从下往上卷紧","5️⃣ 卷好的衣服整齐排列"]
tf=tb(s,0.5,1.9,4.3,2.7,rsteps[0],sz=14,b=True,c=DARK)
for r in rsteps[1:]:ap(tf,"",sz=5);ap(tf,r,sz=14,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.25),Inches(4.6),Inches(3.5))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=TEAL;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.25),Inches(4.6),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=TEAL;rh.line.fill.background()
tb(s,5.25,1.32,4.4,0.4,"🔬 比较实验  Compare",sz=14,b=True,c=WHITE)
tb(s,5.25,1.9,4.35,0.6,"两件一样的 T 恤：一件随便塞，一件用卷衣法。",sz=13,b=True,c=DARK)
tf2=tb(s,5.25,2.6,4.35,1.9,"哪种占的空间更小?",sz=13,c=DARK)
for q in ["哪种更容易找到?","拿出来会不会弄乱其他衣服?"]:
    ap(tf2,"",sz=6);ap(tf2,q,sz=13,c=DARK)
tb(s,5.25,4.25,4.4,0.4,"💡 强调：先铺平再折·两边对齐·别太松·大小相近",sz=10,b=True,c=TEAL)
notes(s,"省空间叠T恤(8分钟)：卷衣法5步。教师强调先铺平再折、两边对齐、不要太松、每件卷成相近大小。比较实验：一件随便塞、一件卷，让学生比空间/好找/不弄乱。")
pn(s,n)

# ============================================================
# 24 省空间叠裤子
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"👖 省空间叠裤子  Fold the Pants",TEAL)
tb(s,0.4,0.85,9.2,0.32,"裤子也能卷得又小又整齐!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.6),Inches(3.5))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=TEAL;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.6),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=TEAL;lh.line.fill.background()
tb(s,0.45,1.32,4.4,0.4,"📏 方法五步  Method",sz=14,b=True,c=WHITE)
psteps=["1️⃣ 把裤子平铺","2️⃣ 两条裤腿对齐","3️⃣ 从裤脚向上折或卷","4️⃣ 调整成整齐小长方形","5️⃣ 和其他裤子放同一区"]
tf=tb(s,0.5,1.9,4.3,2.7,psteps[0],sz=14,b=True,c=DARK)
for p in psteps[1:]:ap(tf,"",sz=5);ap(tf,p,sz=14,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.25),Inches(4.6),Inches(3.5))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=CORAL;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.25),Inches(4.6),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=CORAL;rh.line.fill.background()
tb(s,5.25,1.32,4.4,0.4,"⚠️ 小提醒  Reminders",sz=14,b=True,c=WHITE)
rem=["👖 先把口袋里的东西拿出来","🔗 拉链和扣子先整理好","📐 裤腿要对齐","🧥 厚裤子不要卷得太大"]
tf2=tb(s,5.3,1.95,4.3,2.5,rem[0],sz=14,b=True,c=DARK)
for r in rem[1:]:ap(tf2,"",sz=7);ap(tf2,r,sz=14,b=True,c=DARK)
notes(s,"省空间叠裤子(6分钟)：平铺→裤腿对齐→从裤脚向上折或卷→整齐小长方形→同区。提醒：掏口袋、整理拉链扣子、裤腿对齐、厚裤别卷太大。")
pn(s,n)

# ============================================================
# 25 行李箱摆放挑战
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🧳 行李箱摆放挑战  Pack It Right!",NAVY)
tb(s,0.4,0.85,9.2,0.32,"东西都叠好了 — 怎样放进去最聪明?",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.7),Inches(3.5))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=NAVY;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.7),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=NAVY;lh.line.fill.background()
tb(s,0.45,1.32,4.5,0.4,"📥 摆放顺序  Order",sz=14,b=True,c=WHITE)
order=["1️⃣ 鞋子放底部或侧边","2️⃣ 卷好的衣服整齐排列","3️⃣ 袜子塞进鞋里或小袋","4️⃣ 洗漱用品放防水袋","5️⃣ 小物品放进收纳袋","6️⃣ 常用的放最上面·填满边角"]
tf=tb(s,0.5,1.88,4.4,2.8,order[0],sz=13,b=True,c=DARK)
for o in order[1:]:ap(tf,"",sz=4);ap(tf,o,sz=13,b=True,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.25),Inches(4.5),Inches(3.5))
rp.fill.solid();rp.fill.fore_color.rgb=WARM;rp.line.color.rgb=CORAL;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.25),Inches(4.5),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=CORAL;rh.line.fill.background()
tb(s,5.35,1.32,4.2,0.4,"🤔 想一想  Think",sz=14,b=True,c=WHITE)
tf2=tb(s,5.35,1.9,4.25,2.7,"· 哪些东西能填进小空隙?",sz=13,c=DARK)
for q in ["· 为什么洗漱用品要放防水袋?","· 为什么鞋子不能压干净衣服?","· 回家怎样快速找到衣服?"]:
    ap(tf2,"",sz=6);ap(tf2,q,sz=13,c=DARK)
notes(s,"行李箱摆放挑战(5分钟)：老师示范摆放顺序——鞋子底部/侧边、卷好衣服排列、袜子塞鞋内、洗漱防水袋、小物收纳袋、常用最上、填边角。提问4个。")
pn(s,n)

# ============================================================
# 26 下午任务预告 (第六部分)
# ============================================================
s=ns();n+=1;bg(s,NAVY)
tb(s,0.5,0.55,9,0.7,"📣 下午任务预告  This Afternoon",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.6,1.35,8.8,0.7,"「上午我们学了整理书包、选收纳工具、整理行李箱。\n下午，你们将成为真正的收纳设计师!」",sz=16,b=True,c=SUNNY,a=PP_ALIGN.CENTER)
box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(2.75),Inches(8.0),Inches(1.2))
box.fill.solid();box.fill.fore_color.rgb=WHITE;box.line.color.rgb=CORAL;box.line.width=Pt(2.5)
tb(s,1.2,2.85,7.6,0.4,"👀 上午请特别观察:",sz=14,b=True,c=CORAL)
tb(s,1.2,3.24,7.6,0.65,"哪种整理方法最方便?哪种工具最有用?常用的放哪?\n怎样让别人也看懂你的整理?",sz=13,b=True,c=DARK)
tb(s,0.5,4.2,9,0.35,"📣 老师:「整理不仅要——」 学生:「看起来整齐!」",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,4.58,9,0.35,"老师:「还要——」 学生:「容易找到、容易放回!」",sz=14,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,4.96,9,0.3,"老师:「下午我们要——」 学生:「挑战收纳设计!」",sz=13,b=True,c=SUNNY,a=PP_ALIGN.CENTER)
notes(s,"第六部分(5分钟)：预告下午小组收纳设计挑战，给学生上午观察任务，带结束口号。")
pn(s,n)

# ============================================================
# 27 SESSION 2 DIVIDER — language
# ============================================================
div("Part 6 · 下午","语言目标 + 小组收纳设计挑战\n👀 我会认 5 词  ✍️ 我会写 3 词  🛠️ 收纳挑战",CORAL,"📖")
n+=1

# ============================================================
# 28 QUICK REVIEW — 五步整理法
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🔄 快速复习  Quick Review",CORAL)
tb(s,0.4,0.85,9,0.3,"整理书包五步，你还记得吗?(口头排一排)",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
review=[("🙌","拿出来",TEAL),("🤔","做决定",GOLD),("🗂️","分类",BLUE),("🎒","认识区域",PURPLE),("📥","放回去",GREEN_OK)]
for i,(em,cn,cl) in enumerate(review):
    x=0.35+i*1.9
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.5),Inches(1.7),Inches(2.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.72,1.6,0.8,em,sz=34,a=PP_ALIGN.CENTER)
    tb(s,x,2.62,1.7,0.5,f"{i+1}. {cn}",sz=14,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x,3.2,1.7,0.4,"?",sz=24,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    if i<4:tb(s,x+1.6,1.95,0.4,0.5,"→",sz=22,b=True,c=CORAL)
tb(s,0.4,4.3,9.2,0.4,"💬 先 ___ ，再 ___ ，然后 ___ ，接着 ___ ，最后 ___ 。",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 29-33 我会认 word cards
# ============================================================
def word_card_read(w,py,en,sent,img):
    global n
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",CORAL)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    tb(s,0.5,1.1,4.3,1.4,w,sz=72,b=True,c=TEAL,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.4,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.85,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=CORAL,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=CORAL;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句",sz=16,b=True,c=CORAL)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=22,b=True,c=DARK)
    n+=1;pn(s,n);return s
read_words=[
    ("混乱","hùn luàn","messy","书包太混乱，找不到作业。","📷 混乱"),
    ("分类","fēn lèi","sort","我把东西分类放好。","📷 分类"),
    ("放回","fàng huí","put back","用完要放回原处。","📷 放回"),
    ("收拾","shōu shí","tidy up","玩完玩具要收拾。","📷 收拾"),
    ("书包","shū bāo","backpack","我的书包很整齐。","📷 书包"),
]
for w,py,en,sent,img in read_words:
    word_card_read(w,py,en,sent,img)

# ============================================================
# 34 WORD GAMES
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🎮 练一练  Word Games (选一个玩！)",CORAL)
games=[
    ("1️⃣","拍苍蝇\nFly Swatter","把字卡贴在\n白板上，老师\n说词语，学生拍！",WARM),
    ("2️⃣","举牌游戏\nShow Me","每人 5 张字卡\n老师说词语\n举正确的卡",WARM),
    ("3️⃣","抢椅子\nMusical Chairs","椅子上放字卡\n音乐停，读出词",MINT),
    ("4️⃣","传话筒\nPass the Mic","传球，停下的人\n读字卡并造句",RGBColor(0xE3,0xF2,0xFD)),
]
for i,(num,nm,desc,bgc) in enumerate(games):
    x=0.3+i*2.4
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.9),Inches(2.2),Inches(4.2))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.fill.background()
    tb(s,x+0.1,1.0,2.0,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    ls=nm.split('\n')
    tf=tb(s,x+0.1,1.45,2.0,0.85,ls[0],sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls[1:]:ap(tf,l,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls2=desc.split('\n')
    tf2=tb(s,x+0.15,2.5,1.9,1.5,ls2[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    for l in ls2[1:]:ap(tf2,l,sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,4.75,2.0,0.3,"低 prep ✅",sz=11,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 35-37 我会写 cards
# ============================================================
def word_card_write(w,py,en,img):
    global n
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",TEAL)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=TEAL;sh.line.width=Pt(3)
    tb(s,0.5,1.05,4.3,1.2,w,sz=72,b=True,c=TEAL,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.2,4.3,0.4,f"{py}  {en}",sz=20,c=GRAY,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.0,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.3),Inches(5.0),Inches(1.8))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.fill.background()
    tb(s,0.6,3.4,4.6,0.4,"📝 笔顺 Stroke Order",sz=16,b=True,c=TEAL)
    ib(s,0.6,3.9,4.6,1.0,"📷 插入笔顺图片")
    tf=tb(s,5.8,3.4,3.8,0.4,"练习步骤 Practice:",sz=14,b=True,c=TEAL)
    ap(tf,"1. 空中写 Air Write",sz=13,c=DARK)
    ap(tf,"2. 手心写 Palm Write",sz=13,c=DARK)
    ap(tf,"3. 纸上写 3 times",sz=13,c=DARK)
    n+=1;pn(s,n);return s
write_words=[
    ("分类","fēn lèi","sort","📷 分类"),
    ("书包","shū bāo","backpack","📷 书包"),
    ("收拾","shōu shí","tidy up","📷 收拾"),
]
for w,py,en,img in write_words:
    word_card_write(w,py,en,img)

# ============================================================
# 38 我会说 sentence frames
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"💬 我会说  Sentence Frames (K · G1–3)",CORAL)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.55),Inches(4.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=CORAL;sh.line.width=Pt(2.5)
pb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.45),Inches(1.05),Inches(1.6),Inches(0.4))
pb.fill.solid();pb.fill.fore_color.rgb=CORAL;pb.line.fill.background()
tb(s,0.55,1.10,1.5,0.35,"K (TK-K)",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
k_frames=[("我把 ____ 放回 ____ 。","I put ___ back in ___."),
          ("这是 ____ 的家。","This is ___'s home."),
          ("我会分类。","I can sort."),
          ("我会收拾。","I can tidy up."),
          ("找到了！","Found it!")]
for i,(cn,en) in enumerate(k_frames):
    y=1.55+i*0.62
    tb(s,0.5,y,4.3,0.4,f"·  {cn}",sz=17,b=True,c=TEAL)
    tb(s,0.5,y+0.32,4.3,0.25,en,sz=9,c=GRAY)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.05),Inches(0.95),Inches(4.65),Inches(4.0))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=TEAL;sh2.line.width=Pt(2.5)
pb2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.20),Inches(1.05),Inches(1.7),Inches(0.4))
pb2.fill.solid();pb2.fill.fore_color.rgb=TEAL;pb2.line.fill.background()
tb(s,5.30,1.10,1.6,0.35,"G1 - G3",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
g_frames=[("我把 ___ 和 ___ 放一起，因为 ___ 。","I put ___ and ___ together because ___."),
          ("___ 应该放在 ___ ，因为 ___ 。","___ goes in ___ because ___."),
          ("我们给它贴上「___」标签。","We label it ___."),
          ("这样整理后，可以更快地 ___ 。","Now we can ___ faster.")]
for i,(cn,en) in enumerate(g_frames):
    y=1.65+i*0.78
    tb(s,5.25,y,4.4,0.45,f"·  {cn}",sz=14,b=True,c=TEAL)
    tb(s,5.25,y+0.42,4.4,0.25,en,sz=9,c=GRAY)
tb(s,0.4,5.1,9.2,0.3,"💡 把这张截屏打印，贴在每张桌上 — 学生整堂课参考。",sz=11,b=True,c=NAVY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 39 小组收纳设计挑战 — task + materials + steps
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🛠️ 小组收纳设计挑战  Group Design Challenge",NAVY)
tb(s,0.4,0.85,9.2,0.3,"每组一堆混乱物品 — 你们是收纳设计师!",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.4),Inches(3.6))
lp.fill.solid();lp.fill.fore_color.rgb=WHITE;lp.line.color.rgb=TEAL;lp.line.width=Pt(2.5)
lh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(1.25),Inches(4.4),Inches(0.5))
lh.fill.solid();lh.fill.fore_color.rgb=TEAL;lh.line.fill.background()
tb(s,0.45,1.32,4.2,0.4,"🧺 每组材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.5,1.9,4.1,2.8,"📚 书 · ✏️ 文具 · 📄 纸张",sz=13,c=DARK)
for m in ["🥤 水杯 · 🧸 小玩具 · 🧥 衣服","📦 空盒子 · 🧺 篮子","📁 文件夹 · 🏷️ 标签纸"]:
    ap(tf,"",sz=6);ap(tf,m,sz=13,c=DARK)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.25),Inches(4.8),Inches(3.6))
rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=CORAL;rp.line.width=Pt(2.5)
rh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(1.25),Inches(4.8),Inches(0.5))
rh.fill.solid();rh.fill.fore_color.rgb=CORAL;rh.line.fill.background()
tb(s,5.05,1.32,4.6,0.4,"🎯 小组任务  The Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.05,1.9,4.5,2.8,"1. 全部拿出来 → 分类",sz=12,b=True,c=DARK)
for t in ["2. 决定：保留 / 回收 / 移走","3. 选合适的收纳工具","4. 每类安排固定位置 + 做标签","5. 让别组 10 秒内找到指定物品","6. 根据测试结果修改方案"]:
    ap(tf2,"",sz=4);ap(tf2,t,sz=12,b=True,c=DARK)
notes(s,"下午小组收纳设计挑战：每组一堆混乱物品(书/文具/纸/水杯/玩具/衣服/空盒/篮子/文件夹/标签)。步骤：拿出来→分类→决定保留回收移走→选工具→固定位置+标签→别组10秒找到→修改。")
pn(s,n)

# ============================================================
# 40 评价标准 + 思考
# ============================================================
s=ns();n+=1;bg(s,CREAM);hb(s,"🏆 评价标准  How We Judge",TEAL)
tb(s,0.4,0.85,9.2,0.3,"下午不只看「整不整齐」，还要看这些:",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
crit=["🗂️ 分类是否清楚","📍 位置是否合理","🔎 是否容易找到","✋ 是否方便拿取",
      "🏠 是否容易放回","🏷️ 标签是否有帮助","💧 水杯等是否摆放安全"]
for i,ct in enumerate(crit):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.25+row*0.62
    w=4.55 if not(i==6) else 4.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.52))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=TEAL;sh.line.width=Pt(1.5)
    tb(s,x+0.15,y+0.08,4.3,0.35,ct,sz=14,b=True,c=DARK)
gold=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.75),Inches(9.4),Inches(1.1))
gold.fill.solid();gold.fill.fore_color.rgb=WARM;gold.line.color.rgb=CORAL;gold.line.width=Pt(2.5)
tb(s,0.5,3.83,9.0,0.35,"💬 结束口号  Chant:",sz=13,b=True,c=CORAL)
tb(s,0.5,4.18,9.2,0.35,"老师:「整理不仅要——」学生:「看起来整齐!」",sz=13,b=True,c=DARK)
tb(s,0.5,4.5,9.2,0.35,"老师:「还要——」学生:「容易找到、容易放回!」",sz=13,b=True,c=DARK)
notes(s,"评价标准7项：分类清楚/位置合理/容易找到/方便拿取/容易放回/标签有帮助/水杯安全。用结束口号收尾。")
pn(s,n)

# ============================================================
# 41 Day 2 Badge
# ============================================================
s=ns();n+=1;bg(s,CREAM)
tb(s,0.5,0.3,9,0.7,"🎖️ Day 2 整理达人徽章  Master Badge",sz=24,b=True,c=TEAL,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(1.05),Inches(3),Inches(3))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=TEAL;sh.line.width=Pt(5)
tf=tb(s,3.6,1.28,2.8,0.4,"DAY 2",sz=18,b=True,c=CORAL,a=PP_ALIGN.CENTER)
ap(tf,"📦",sz=40,a=PP_ALIGN.CENTER)
ap(tf,"整理和收纳",sz=19,b=True,c=TEAL,a=PP_ALIGN.CENTER)
ap(tf,"✓ COMPLETED",sz=13,b=True,c=GREEN_OK,a=PP_ALIGN.CENTER)
ap(tf,"🎒🧰🧳✅",sz=16,a=PP_ALIGN.CENTER)
sb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.3),Inches(4.12),Inches(7.4),Inches(0.6))
sb.fill.solid();sb.fill.fore_color.rgb=SUNNY;sb.line.color.rgb=CORAL;sb.line.width=Pt(2.5)
tb(s,1.3,4.15,7.4,0.55,"⭐  ⭐  ⭐  ⭐  ⭐  ⭐",sz=30,b=True,c=CORAL,a=PP_ALIGN.CENTER)
tb(s,1,4.8,8,0.4,"📣 老师:「用完东西——」 学生:「送它回家!」🎉",sz=15,b=True,c=TEAL,a=PP_ALIGN.CENTER)
tb(s,1,5.15,8,0.3,"五步整理书包 · 收纳工具 · 行李箱叠衣法 · 5 个词",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
pn(s,n)

# ============================================================
# 42 Tomorrow preview
# ============================================================
s=ns();n+=1;bg(s,TEAL)
tb(s,0.5,0.9,9,0.8,"🌟 明天见！  See You Tomorrow!",sz=32,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tf=tb(s,1.5,2.2,7,2.5,"Day 3 — 待定 (To Be Continued)",sz=26,b=True,c=SUNNY,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"🏠 继续做小小生活家！",sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
ap(tf,"Keep being a little homemaker!",sz=14,c=WARM,a=PP_ALIGN.CENTER)
ap(tf,"",sz=10)
ap(tf,"明天见，整理达人！",sz=15,c=WARM,a=PP_ALIGN.CENTER)
pn(s,n)

OUT='/Users/huanli/projects/courseppt/Chinese/小小生活家/day2_organizing.pptx'
prs.save(OUT);print(f"Created {n} slides → {OUT}")
