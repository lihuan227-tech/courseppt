# create_day3_ocean.py — Zero Waste Day 3 · 海洋小卫士 (Water + Plastic)
# Follows the wilderness day1_nature reference structure (看→想→判→做), ocean-themed.
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import _helpers as H
from _helpers import (make_presentation, ns, tb, bg, hb, pn, panel, panel_head, div,
                      cover, learning_goals, photo_slot, teacher_box, sentence_frame_bar,
                      video_slide, vocab_recognize, vocab_write, share_close,
                      mission_intro_slide, observe_think_slide, judge_ab_slide,
                      reveal_ab_slide, cardgrid_slide, guess_slide,
                      CREAM, WHITE, DARK, GRAY, LGRAY, INK, STAR, WARM, IMGBG,
                      OK, ALERT, FIRE_ORANGE, EARTH_BROWN)

# ----- Ocean palette -----
def _rgb(h): return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
OCEAN = _rgb("0B5C8C")
TEAL  = _rgb("006970")
AQUA  = _rgb("2A9D8F")
SKY   = _rgb("42A5F5")
CORAL = ALERT
TINTS = list(H.OCEAN_TINTS)

# ---------------- content data ----------------
# item: em, cn(spaced), short, en, tint, victim, senses[3], dangers[4], rows[3], do, swap, media
ITEMS = [
  dict(em="🛍️", cn="塑 料 袋", short="塑料袋", en="Plastic Bag", tint=0, victim="🐢",
       senses=[("👀", "薄 薄 的, 会 飞"), ("🌊", "在 水 里 像 水 母"), ("🗑️", "用 一 次 就 扔?")],
       dangers=[("🐢", "海龟?"), ("🐟", "鱼?"), ("🐦", "海鸟?"), ("🌊", "漂到海里?")],
       rows=[("去 买 东 西", "要 塑 料 袋", "带 布 袋", "B"),
             ("袋 子 破 了", "再 拿 一 个", "重 复 用", "B"),
             ("买 菜", "很 多 小 袋", "一 个 布 袋", "B")],
       do="假 装 把 布 袋 挂 手 上, 去 买 东 西!", swap="🛍️ 换 成: 布 袋 / 可 重 复 袋", media=None),
  dict(em="🥤", cn="吸 管", short="吸管", en="Straw", tint=1, victim="🐢",
       senses=[("👀", "细 细 的, 一 次 性"), ("🌊", "会 掉 进 海 里"), ("🐢", "卡 在 海 龟 鼻 子 里?")],
       dangers=[("🐢", "海龟?"), ("🐬", "海豚?"), ("🐟", "鱼?"), ("🗑️", "变垃圾?")],
       rows=[("喝 果 汁", "用 塑 料 吸 管", "不 用, 直 接 喝", "B"),
             ("店 里 给 吸 管", "收 下", "说 不 用", "B"),
             ("想 要 吸 管", "塑 料 的", "纸 的 / 钢 的", "B")],
       do="做 喝 水 动 作 — 不 用 吸 管!", swap="🚫 不 用 · 钢 吸 管 / 纸 吸 管",
       media=("看 海 龟 与 吸 管 视 频", "https://www.youtube.com/results?search_query=sea+turtle+straw")),
  dict(em="🧴", cn="塑 料 瓶", short="塑料瓶", en="Plastic Bottle", tint=2, victim="🐦",
       senses=[("👀", "透 明, 很 多"), ("🌊", "会 碎 成 小 塑 料"), ("🗑️", "喝 完 就 扔?")],
       dangers=[("🐦", "海鸟?"), ("🐟", "鱼?"), ("🐳", "鲸鱼?"), ("🌊", "海洋垃圾?")],
       rows=[("口 渴 了", "买 瓶 装 水", "带 水 壶", "B"),
             ("水 喝 完 了", "再 买 一 瓶", "装 满 水 壶", "B"),
             ("出 门", "买 饮 料 瓶", "自 己 水 壶", "B")],
       do="举 起 水 壶, 咕 嘟 咕 嘟 喝 水!", swap="💧 换 成: 水 壶 (装 满 再 用)", media=None),
  dict(em="🍴", cn="一 次 性 餐 具", short="餐具", en="Disposable Utensils", tint=3, victim="🐟",
       senses=[("👀", "塑 料 叉 / 勺"), ("🌊", "会 断 成 小 片"), ("🗑️", "吃 完 就 扔?")],
       dangers=[("🐟", "鱼?"), ("🦀", "螃蟹?"), ("🐢", "海龟?"), ("🗑️", "变垃圾?")],
       rows=[("吃 饭", "一 次 性 叉", "自 带 勺 子", "B"),
             ("外 卖 给 餐 具", "都 收 下", "说 不 用", "B"),
             ("野 餐", "塑 料 餐 具", "自 己 的", "B")],
       do="假 装 从 书 包 拿 出 勺 子 吃 饭!", swap="🍴 换 成: 自 带 餐 具", media=None),
  dict(em="📦", cn="塑 料 包 装", short="包装", en="Plastic Packaging", tint=4, victim="🦭",
       senses=[("👀", "一 层 层 包 装"), ("🌊", "会 缠 住 动 物"), ("🗑️", "拆 完 就 扔?")],
       dangers=[("🦭", "海豹?"), ("🐢", "海龟?"), ("🐟", "鱼?"), ("🌊", "漂到海里?")],
       rows=[("买 零 食", "很 多 小 包 装", "大 包 装", "B"),
             ("买 水 果", "塑 料 盒 装", "散 装 自 己 装", "B"),
             ("包 装 绳", "套 着 扔", "剪 开 再 扔", "B")],
       do="用 手 比 一 比 — 剪 开 包 装 再 扔!", swap="📦 换 成: 少 包 装 / 散 装", media=None),
  dict(em="🥤", cn="一 次 性 杯 子", short="杯子", en="Disposable Cup", tint=5, victim="🌊",
       senses=[("👀", "纸 / 塑 料 杯 + 盖"), ("🌊", "会 变 海 洋 垃 圾"), ("🗑️", "喝 完 就 扔?")],
       dangers=[("🌊", "海洋垃圾?"), ("🐦", "海鸟?"), ("🐟", "鱼?"), ("🐢", "海龟?")],
       rows=[("买 饮 料", "一 次 性 杯", "自 己 的 杯", "B"),
             ("喝 水", "拿 纸 杯", "用 自 己 杯", "B"),
             ("咖 啡 店", "要 一 次 性 杯", "带 随 行 杯", "B")],
       do="举 起 自 己 的 杯 子, 干 杯!", swap="🥤 换 成: 自 带 杯 子 / 随 行 杯", media=None),
]

VOCAB_RECOGNIZE = [
  ("💧", "水", "shuǐ", "water", "水 是 生 命 之 源, 每 天 都 要 喝 水。", "Water is the source of life.", "💧 一 杯 干 净 的 水"),
  ("🌊", "海 洋", "hǎi yáng", "ocean", "海 洋 里 有 很 多 动 物。", "The ocean has many animals.", "🌊 蓝 色 的 大 海"),
  ("🥤", "塑 料", "sù liào", "plastic", "塑 料 用 一 次 就 扔, 很 浪 费。", "Single-use plastic is wasteful.", "🥤 塑 料 瓶 和 塑 料 袋"),
  ("🛢️", "污 染", "wū rǎn", "pollution", "塑 料 让 海 洋 污 染。", "Plastic pollutes the ocean.", "🛢️ 海 面 上 的 塑 料 垃 圾"),
  ("🛡️", "保 护", "bǎo hù", "protect", "我 们 要 保 护 海 洋。", "We protect the ocean.", "🛡️ 孩 子 在 捡 垃 圾"),
  ("🚰", "节 约", "jié yuē", "save", "节 约 用 水, 关 好 水 龙 头。", "Save water — turn off the tap.", "🚰 关 水 龙 头 的 手"),
]

VOCAB_WRITE = [
  ("水", [("水", "shuǐ", "4 笔", "三 点 加 一 竖 钩")], "water"),
  ("保 护", [("保", "bǎo", "9 笔", "亻 加 呆"), ("护", "hù", "7 笔", "扌 加 户")], "protect"),
  ("节 约", [("节", "jié", "5 笔", "艹 加 卩"), ("约", "yuē", "6 笔", "纟 加 勺")], "save"),
]


def build():
    prs = make_presentation()
    n = [0]
    def page(s):
        n[0] += 1; return s  # numbering stamped in a final pass so pn draws on top

    # ---------- Open (1-3) ----------
    page(cover(prs, 3, "水 与 塑 料 · 海 洋 小 卫 士",
               "Water & Plastic — Little Ocean Guardians",
               "💧  🌊  🐢  🥤  🛍️  ♻️", OCEAN,
               "我 们 怎 样 节 约 用 水、减 少 塑 料, 保 护 海 洋?",
               "How can we save water and cut plastic to protect the ocean?"))

    s = page(ns(prs)); bg(s, CREAM); hb(s, "⏰ 今 日 时 间 安 排  Today's Schedule", OCEAN)
    for i, (lab, tm, desc, cl) in enumerate([
            ("Session 1 · 上午", "10:00–11:45", "水 很 重 要 + 塑 料 污 染 (故 事/视 频)", OCEAN),
            ("Session 2 · 下午", "2:00–2:45", "复 习 + 语 言 (我 会 认 6 / 我 会 写 3)", TEAL),
            ("Session 3 · 下午", "3:00–4:30", "练 习 册 + 动 手 (海 洋 拼 贴 / 承 诺)", AQUA)]):
        y = 1.05 + i*1.42
        p = panel(s, 0.5, y, 9.0, 1.25, cl, fill=cl, lw=0); p.line.fill.background()
        tb(s, 0.8, y+0.22, 5.3, 0.5, lab, sz=22, b=True, c=WHITE)
        tb(s, 0.8, y+0.78, 5.3, 0.35, tm, sz=13, c=STAR)
        tb(s, 5.7, y+0.40, 3.6, 0.8, desc, sz=13, b=True, c=WHITE)

    page(learning_goals(prs, OCEAN, [
        ("1", "理 解 水 资 源 的 重 要 性", "Understand why water matters", OCEAN),
        ("2", "认 识 塑 料 对 海 洋 动 物 的 影 响", "Plastic harms ocean animals", TEAL),
        ("3", "理 解 一 次 性 塑 料 的 问 题", "The problem with single-use plastic", AQUA),
        ("4", "提 出 减 少 塑 料 的 方 法", "Propose ways to reduce plastic", CORAL)]))

    # ---------- Session 1 (4-34) ----------
    page(div(prs, "Session 1 · 上 午", "水 很 重 要 + 塑 料 污 染", OCEAN, "🌊"))

    page(mission_intro_slide(prs, OCEAN, "🧭 你 是 海 洋 小 卫 士!  Little Ocean Guardian!",
        [(it["em"], it["short"], it["en"]) for it in ITEMS],
        "我 是 海 洋 小 卫 士, 我 要 减 少 ______。", "I'm an Ocean Guardian. I'll cut down on ___."))

    s = page(ns(prs)); bg(s, CREAM); hb(s, "💬 说 一 说 · 水  Let's Talk About Water", OCEAN)
    for i, (q, en) in enumerate([("你 在 哪 里 见 过 水? (河、海、雨、水 龙 头)", "Where have you seen water?"),
                                 ("我 们 为 什 么 需 要 水?", "Why do we need water?")]):
        y = 1.0 + i*1.0; panel(s, 0.5, y, 9.0, 0.85, OCEAN, fill=WHITE, lw=2.5)
        tb(s, 0.7, y+0.14, 8.6, 0.4, f"{i+1}. {q}", sz=15, b=True, c=DARK)
        tb(s, 0.7, y+0.54, 8.6, 0.3, en, sz=10, c=GRAY)
    sentence_frame_bar(s, 3.05, "我 们 需 要 水 来 ______。", "We need water to ___.", OCEAN)
    teacher_box(s, 0.5, 3.75, 9.0, "引 导 说 出: 喝、洗、种 植 物、动 物 也 要 水", "学 生 抢 答", 4, OCEAN)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "💧 水 很 重 要  Water Matters", OCEAN)
    tb(s, 0.4, 0.82, 9.2, 0.3, "我 们 每 天 都 要 用 水!  We use water every day!", sz=13, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    for i, (em, cn) in enumerate([("🥤", "喝 水"), ("🚿", "洗 澡"), ("🌱", "种 植 物"), ("🐟", "动 物 也 要 水")]):
        r, c = divmod(i, 2); x = 0.9 + c*4.3; y = 1.3 + r*1.65
        panel(s, x, y, 3.9, 1.45, OCEAN, fill=WHITE, lw=2.5)
        tb(s, x, y+0.16, 3.9, 0.68, em, sz=42, a=PP_ALIGN.CENTER)
        tb(s, x, y+0.92, 3.9, 0.4, cn, sz=18, b=True, c=OCEAN, a=PP_ALIGN.CENTER)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🌍 地 球 的 水 很 珍 贵 · 节 约 用 水  Save Water", TEAL)
    panel(s, 0.4, 1.0, 4.5, 3.4, TEAL, fill=WARM, lw=2.5)
    tb(s, 0.5, 1.35, 4.3, 1.2, "🌏", sz=88, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 2.95, 4.3, 0.5, "能 喝 的 淡 水 很 少!", sz=17, b=True, c=TEAL, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 3.55, 4.3, 0.5, "Only a little water is drinkable.", sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    panel(s, 5.1, 1.0, 4.5, 3.4, TEAL, fill=WHITE, lw=2.5); panel_head(s, 5.1, 1.0, 4.5, TEAL, "✅ 节 约 用 水 小 妙 招", sz=13)
    for i, (em, cn) in enumerate([("🚰", "关 好 水 龙 头"), ("🪥", "刷 牙 时 关 水"),
                                  ("♻️", "洗 菜 水 浇 花"), ("🚿", "洗 澡 快 一 点")]):
        y = 1.65 + i*0.63; tb(s, 5.25, y, 0.5, 0.5, em, sz=22); tb(s, 5.85, y+0.06, 3.6, 0.4, cn, sz=13, b=True, c=DARK)

    # ---- 绘本 Somebody Swallowed Stanley: 塑料袋 → 海洋 → 动物 → 我们 ----
    s = page(ns(prs)); bg(s, CREAM)
    hb(s, "📖 绘 本 时 间 · 谁 吞 掉 了 斯 坦 利?  Storybook Time", OCEAN)
    panel(s, 0.40, 0.95, 4.40, 3.95, OCEAN, fill=INK, lw=3)
    tb(s, 0.40, 1.15, 4.40, 1.05, "🛍️", sz=78, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 2.28, 4.40, 0.45, "谁 吞 掉 了 斯 坦 利?", sz=20, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 2.75, 4.40, 0.30, "Somebody Swallowed Stanley", sz=12, c=WARM, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 3.05, 4.40, 0.26, "中 文 绘 本 · 拯 救 海 洋", sz=10, c=LGRAY, a=PP_ALIGN.CENTER)
    pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.20), Inches(3.50), Inches(2.80), Inches(0.60))
    pb.fill.solid(); pb.fill.fore_color.rgb = FIRE_ORANGE; pb.line.fill.background()
    tb(s, 1.20, 3.60, 2.80, 0.42, "▶️  读 绘 本  Read Aloud", sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    _url = "https://www.youtube.com/watch?v=LHY2XcCiBSw"
    pb.click_action.hyperlink.address = _url
    tb(s, 0.40, 4.32, 4.40, 0.28, "🛍️ 斯 坦 利 是 一 个 塑 料 袋", sz=10, c=STAR, a=PP_ALIGN.CENTER)
    panel(s, 5.10, 0.95, 4.50, 3.95, FIRE_ORANGE, fill=WHITE, lw=3)
    panel_head(s, 5.10, 0.95, 4.50, FIRE_ORANGE, "🎬 读 之 前 — 想 一 想", sz=13)
    for i, (em, cn, en) in enumerate([
            ("🛍️", "斯 坦 利 是 什 么?", "What is Stanley?"),
            ("🌊", "它 掉 进 海 里 会 怎 样?", "What happens in the sea?"),
            ("🐢", "谁 会 把 它 当 成 食 物?", "Who thinks it's food?")]):
        y = 1.60 + i*1.00
        tb(s, 5.25, y, 0.55, 0.55, em, sz=28)
        tb(s, 5.85, y+0.08, 3.60, 0.38, cn, sz=13, b=True, c=DARK)
        tb(s, 5.85, y+0.48, 3.60, 0.30, en, sz=9, c=GRAY)
    ab = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.45))
    ab.fill.solid(); ab.fill.fore_color.rgb = OCEAN; ab.line.fill.background()
    tb(s, 0.55, 5.07, 9.0, 0.32, "📖 一 边 读, 一 边 想: 塑 料 袋 给 海 洋 带 来 什 么 麻 烦?",
       sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)

    # ---- 绘本讨论 + 对人的危害 (food chain: plastic → fish → us) ----
    s = page(ns(prs)); bg(s, CREAM)
    hb(s, "💬 绘 本 讨 论 · 塑 料 也 会 伤 害 我 们  Let's Talk", CORAL)
    panel(s, 0.40, 1.00, 4.45, 3.05, OCEAN, fill=WHITE, lw=2.5)
    panel_head(s, 0.40, 1.00, 4.45, OCEAN, "🤔 想 一 想 · 绘 本 里", sz=13)
    for i, (em, q) in enumerate([
            ("🐢", "斯 坦 利 让 哪 些 动 物 受 伤?"),
            ("🌊", "海 龟 为 什 么 会 吃 塑 料 袋?"),
            ("😊", "最 后 谁 救 了 海 龟?")]):
        y = 1.62 + i*0.76
        tb(s, 0.58, y, 0.55, 0.5, em, sz=22)
        tb(s, 1.15, y+0.06, 3.55, 0.5, q, sz=13, b=True, c=DARK)
    panel(s, 5.10, 1.00, 4.50, 3.05, CORAL, fill=WARM, lw=2.5)
    panel_head(s, 5.10, 1.00, 4.50, CORAL, "🍽️ 塑 料 会 回 到 我 们 身 上!", sz=12)
    for i, (em, cn) in enumerate([
            ("🥤", "塑 料 碎 成 小 片"), ("🐟", "小 鱼 吃 塑 料"),
            ("🐠", "大 鱼 吃 小 鱼"), ("🧑", "我 们 吃 鱼")]):
        y = 1.55 + i*0.52
        tb(s, 5.30, y, 0.55, 0.45, em, sz=20)
        tb(s, 5.92, y+0.04, 3.4, 0.4, cn, sz=13, b=True, c=DARK)
    tb(s, 5.25, 3.66, 4.20, 0.32, "→ 海 洋 垃 圾 也 伤 害 人!", sz=12, b=True, c=CORAL, a=PP_ALIGN.CENTER)
    sentence_frame_bar(s, 4.25, "塑 料 伤 害 ______, 也 伤 害 我 们。", "Plastic hurts ___ and us too.", OCEAN)
    nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.30), Inches(4.92), Inches(9.40), Inches(0.48))
    nb.fill.solid(); nb.fill.fore_color.rgb = WARM; nb.line.color.rgb = CORAL; nb.line.width = Pt(1.5)
    tb(s, 0.45, 5.00, 9.10, 0.32, "👩‍🏫 先 读 绘 本, 再 讨 论 — 强 调: 塑 料 → 鱼 → 我 们  (⏱️ 8 分)",
       sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🥤 什 么 是 一 次 性 塑 料?  Single-Use Plastic", OCEAN)
    tb(s, 0.4, 0.88, 9.2, 0.5, "用 一 次 就 扔 掉 的 塑 料 = 一 次 性 塑 料", sz=20, b=True, c=CORAL, a=PP_ALIGN.CENTER)
    tb(s, 0.4, 1.45, 9.2, 0.35, "Used once, then thrown away.", sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    for i, it in enumerate(ITEMS):
        r, c = divmod(i, 3); x = 0.7 + c*3.0; y = 2.05 + r*1.5
        panel(s, x, y, 2.8, 1.35, TINTS[it["tint"]], fill=WHITE, lw=2)
        tb(s, x, y+0.15, 2.8, 0.6, it["em"], sz=34, a=PP_ALIGN.CENTER)
        tb(s, x, y+0.85, 2.8, 0.4, it["cn"], sz=15, b=True, c=TINTS[it["tint"]], a=PP_ALIGN.CENTER)

    page(cardgrid_slide(prs, OCEAN, "🗺️ 今 天 的 6 种 一 次 性 塑 料  Today's 6 Single-Use Plastics",
        "先 认 一 认, 再 一 个 一 个 来!",
        [(it["em"], it["cn"], it["en"]) for it in ITEMS]))

    for it in ITEMS:
        col = TINTS[it["tint"]]
        hdr = f'{it["em"]} {it["cn"]} · {it["en"]}'
        media_cn, media_url = it["media"] if it["media"] else (None, None)
        page(observe_think_slide(prs, col, f"{hdr}  ①看+②想",
            it["cn"].replace(" ", "") + " 的 真 实 照 片", f'Real photo of {it["en"]}',
            it["senses"], it["dangers"], media_cn, media_url))
        page(judge_ab_slide(prs, col, f'{it["em"]} {it["cn"]} · ⚖️ 你 觉 得 呢?', it["rows"],
            "K-G1: 我 选 A / B。   G2-G3: 我 选 __, 因 为 __。", "I choose A/B. (…because…)"))
        page(reveal_ab_slide(prs, col, f'{it["em"]} {it["cn"]} · 💡 答 案 揭 晓!', it["rows"],
            it["do"] + "   " + it["swap"],
            "____ 不 好, 我 应 该 ____。", "___ is bad; I should ___."))

    # 31 compare table
    s = page(ns(prs)); bg(s, CREAM); hb(s, "📋 6 种 塑 料 对 比  Compare 6 Plastics", OCEAN)
    tb(s, 0.4, 0.80, 9.2, 0.3, "每 种 塑 料 伤 害 谁? 换 成 什 么?", sz=12, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    swaps = ["🛍️", "🚫", "💧", "🍴", "📦", "🥤"]
    rows_data = [("塑 料", [it["em"] for it in ITEMS]),
                 ("动 物", [it["victim"] for it in ITEMS]),
                 ("换 成", swaps)]
    ty = 1.55
    for ri, (lab, vals) in enumerate(rows_data):
        y = ty + ri*1.25
        p = panel(s, 0.35, y, 1.15, 1.05, OCEAN, fill=OCEAN, lw=0); p.line.fill.background()
        tb(s, 0.35, y+0.34, 1.15, 0.4, lab, sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
        for ci, v in enumerate(vals):
            x = 1.60 + ci*1.33; panel(s, x, y, 1.25, 1.05, OCEAN, fill=WHITE, lw=1.5)
            tb(s, x, y+0.30, 1.25, 0.6, v, sz=26, a=PP_ALIGN.CENTER)

    # 32-33 action game
    game = [("🛍️", "买 东 西 不 用 塑 料 袋"), ("🥤", "喝 水 用 自 己 的 水 壶"), ("🐢", "塑 料 扔 进 海 里"),
            ("💧", "刷 牙 一 直 开 着 水"), ("🍴", "自 带 勺 子 吃 饭"), ("🚰", "洗 菜 水 浇 花")]
    ans = ["✅", "✅", "❌", "❌", "✅", "✅"]
    s = page(ns(prs)); bg(s, CREAM); hb(s, "🛡️ 好 vs 不 好!  🤔 你 觉 得 呢?", OCEAN)
    tb(s, 0.4, 0.85, 9.2, 0.4, "老 师 说 一 件 事 — 学 生 做 动 作!  ✅ 好 = 举 手 · ❌ 不 好 = 交 叉", sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    for i, (em, cn) in enumerate(game):
        r, c = divmod(i, 2); x = 0.5 + c*4.6; y = 1.5 + r*1.15
        panel(s, x, y, 4.3, 1.0, OCEAN, fill=WHITE, lw=2); tb(s, x+0.12, y+0.28, 0.6, 0.5, em, sz=26)
        tb(s, x+0.78, y+0.30, 2.9, 0.5, cn, sz=13, b=True, c=DARK); tb(s, x+3.7, y+0.24, 0.5, 0.5, "?", sz=24, b=True, c=OCEAN)
    s = page(ns(prs)); bg(s, CREAM); hb(s, "🛡️ 好 vs 不 好!  💡 答 案 揭 晓", CORAL)
    for i, (em, cn) in enumerate(game):
        r, c = divmod(i, 2); x = 0.5 + c*4.6; y = 1.3 + r*1.13; mk = ans[i]; mc = OK if mk == "✅" else ALERT
        panel(s, x, y, 4.3, 1.0, mc, fill=WHITE, lw=2.5); tb(s, x+0.12, y+0.28, 0.6, 0.5, em, sz=26)
        tb(s, x+0.78, y+0.30, 2.8, 0.5, cn, sz=13, b=True, c=DARK); tb(s, x+3.62, y+0.22, 0.6, 0.55, mk, sz=26, b=True, c=mc)
    sentence_frame_bar(s, 5.0, "____ 不 好, 因 为 ____。", "___ is bad, because ___.", OCEAN)

    page(guess_slide(prs, OCEAN, "🔍 我 演 你 猜  Act & Guess!",
        "老 师 演 一 个 动 作 — 学 生 猜 是 哪 个 塑 料 / 替 代 品!",
        [("🛍️", "挂 布 袋 去 买 东 西"), ("🥤", "咕 嘟 咕 嘟 喝 水 壶"), ("🍴", "从 书 包 拿 勺 子"),
         ("🚫", "摆 手 说 不 用 吸 管"), ("📦", "剪 开 包 装"), ("🥤", "举 杯 干 杯")],
        "我 猜 是 ______!", "I guess it's ___!"))

    # ---------- Session 2 (35-46) ----------
    page(div(prs, "Session 2 · 下 午", "复 习 + 语 言 (我 会 认 6 · 我 会 写 3)", TEAL, "📖"))

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🔁 快 速 复 习 — 塑 料 伤 害 谁?", TEAL)
    tb(s, 0.4, 0.82, 9.2, 0.3, "把 塑 料 和 它 伤 害 的 动 物 连 起 来 (口 头)", sz=12, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    for i, it in enumerate(ITEMS):
        y = 1.25 + i*0.62
        p = panel(s, 1.3, y, 3.2, 0.52, TINTS[it["tint"]], fill=TINTS[it["tint"]], lw=0); p.line.fill.background()
        tb(s, 1.45, y+0.10, 3.0, 0.35, f'{it["em"]} {it["cn"]}', sz=14, b=True, c=WHITE)
        panel(s, 5.5, y, 3.2, 0.52, CORAL, fill=WHITE, lw=2)
        tb(s, 5.65, y+0.10, 3.0, 0.35, f'{it["victim"]}  ？', sz=14, b=True, c=DARK)
    tb(s, 4.55, 2.4, 0.9, 0.6, "❓", sz=28, a=PP_ALIGN.CENTER)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🎮 复 习 游 戏 · Baamboozle", TEAL)
    panel(s, 1.5, 1.1, 7.0, 2.6, TEAL, fill=INK, lw=3); tb(s, 1.5, 1.85, 7.0, 1.0, "🎮", sz=80, a=PP_ALIGN.CENTER)
    btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(3.95), Inches(4.0), Inches(0.6))
    btn.fill.solid(); btn.fill.fore_color.rgb = FIRE_ORANGE; btn.line.fill.background()
    tb(s, 3.0, 4.05, 4.0, 0.4, "▶️ 点 击 开 始  Play", sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    url = "https://www.baamboozle.com/games"  # TODO: replace with the real Day-3 game id
    btn.click_action.hyperlink.address = url
    tb(s, 0.4, 4.75, 9.2, 0.3, "🔗 " + url + "  (老 师: 换 成 真 实 游 戏 链 接)", sz=9, c=GRAY, a=PP_ALIGN.CENTER)

    for (em, cn, py, en, ex, exen, hint) in VOCAB_RECOGNIZE:
        page(vocab_recognize(prs, OCEAN, em, cn, py, en, ex, exen, hint))

    for (phrase, chars, en) in VOCAB_WRITE:
        page(vocab_write(prs, OCEAN, phrase, en, chars))

    # ---------- Session 3 (47-54) ----------
    page(div(prs, "Session 3 · 下 午", "练 习 册 + 动 手 (海 洋 拼 贴 · 承 诺)", AQUA, "🎒"))

    s = page(ns(prs)); bg(s, CREAM); hb(s, "📓 完 成 「海 洋 小 卫 士」练 习 册  Day 3 Booklet", EARTH_BROWN)
    for i, (t, d) in enumerate([("① 圈 污 染", "哪 些 是 一 次 性 塑 料?"), ("② 连 一 连", "塑 料 → 换 成 什 么"),
                                ("③ 我 的 承 诺", "我 是 海 洋 小 卫 士, 我 可 以 ___"), ("④ 描 一 描", "水 · 保 护 · 塑 料")]):
        r, c = divmod(i, 2); x = 0.6 + c*4.5; y = 1.22 + r*1.98
        panel(s, x, y, 4.2, 1.70, AQUA, fill=WHITE, lw=2.5); tb(s, x+0.25, y+0.30, 3.7, 0.5, t, sz=17, b=True, c=AQUA)
        tb(s, x+0.25, y+0.95, 3.7, 0.6, d, sz=12, b=True, c=DARK)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🎨 动 手 时 间!  Hands-On — 2 个 活 动", AQUA)
    for i, (t, d, cl, em) in enumerate([("PROJECT 1 · 海 洋 拼 贴 画", "用 回 收 材 料 拼 一 片 干 净 的 海 洋", OCEAN, "🖼️"),
                                        ("PROJECT 2 · 装 饰 环 保 袋", "装 饰 布 袋 + 写 保 护 海 洋 的 承 诺", AQUA, "🛍️")]):
        x = 0.5 + i*4.7; panel(s, x, 1.1, 4.3, 3.4, cl, fill=WHITE, lw=3)
        tb(s, x+0.2, 1.35, 3.9, 0.8, t, sz=16, b=True, c=cl, a=PP_ALIGN.CENTER)
        tb(s, x+0.2, 2.2, 3.9, 0.9, em, sz=68, a=PP_ALIGN.CENTER)
        tb(s, x+0.2, 3.5, 3.9, 0.8, d, sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🖼️ Project 1: 海 洋 拼 贴 画  Ocean Collage", OCEAN)
    panel_head(s, 0.4, 0.95, 4.5, OCEAN, "🧺 材 料  Materials")
    for i, m in enumerate(["🟦 蓝 纸 (海 水)", "🐟 彩 纸 (鱼)", "♻️ 干 净 塑 料 片", "🍃 树 叶", "🖊️ 胶 水"]):
        tb(s, 0.55, 1.55+i*0.5, 4.2, 0.4, m, sz=13, b=True, c=DARK)
    panel_head(s, 5.1, 0.95, 4.5, FIRE_ORANGE, "👉 做 法  Steps")
    for i, st in enumerate(["1️⃣ 想 一 片 干 净 的 海 洋", "2️⃣ 用 材 料 拼 出 海 和 动 物", "3️⃣ 贴 牢, 晾 一 下", "4️⃣ 介 绍 我 的 作 品"]):
        tb(s, 5.25, 1.55+i*0.55, 4.3, 0.45, st, sz=13, b=True, c=DARK)
    sentence_frame_bar(s, 4.7, "这 是 干 净 的 海 洋, 有 ______。", "This is a clean ocean with ___.", OCEAN)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🖼️ Project 1 · 参 考 作 品  Examples", OCEAN)
    for i in range(6):
        r, c = divmod(i, 3); x = 0.45 + c*3.13; y = 0.98 + r*2.20
        photo_slot(s, x, y, 2.83, 1.90, "海 洋 拼 贴 参 考", "collage example", OCEAN)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🛍️ Project 2: 装 饰 环 保 袋 + 承 诺", AQUA)
    panel_head(s, 0.4, 0.95, 4.5, AQUA, "🧺 材 料 + 做 法  Materials + Steps")
    for i, st in enumerate(["🛍️ 一 个 布 袋", "🖍️ 彩 笔 / 贴 纸", "1️⃣ 画 海 洋 动 物", "2️⃣ 写 一 句 承 诺", "3️⃣ 贴 到 承 诺 墙"]):
        tb(s, 0.55, 1.55+i*0.55, 4.2, 0.45, st, sz=13, b=True, c=DARK)
    panel(s, 5.1, 0.95, 4.5, 3.3, AQUA, fill=WARM, lw=2.5); tb(s, 5.1, 1.55, 4.5, 1.4, "🛍️🌊", sz=68, a=PP_ALIGN.CENTER)
    tb(s, 5.2, 3.2, 4.3, 0.8, "我 的 布 袋, 我 天 天 用!", sz=15, b=True, c=AQUA, a=PP_ALIGN.CENTER)
    sentence_frame_bar(s, 4.7, "我 是 海 洋 小 卫 士, 我 可 以 ______。", "I'm an Ocean Guardian. I can ___.", AQUA)

    s = page(ns(prs)); bg(s, CREAM); hb(s, "🏅 Day 3 海 洋 小 卫 士 徽 章  Ocean Guardian Badge", OCEAN)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(0.95), Inches(3.6), Inches(3.0))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.color.rgb = OCEAN; c.line.width = Pt(4)
    tb(s, 3.2, 1.15, 3.6, 0.4, "DAY 3", sz=16, b=True, c=CORAL, a=PP_ALIGN.CENTER)
    tb(s, 3.2, 1.55, 3.6, 0.9, "🌊🐢", sz=52, a=PP_ALIGN.CENTER)
    tb(s, 3.2, 2.62, 3.6, 0.5, "海 洋 小 卫 士", sz=20, b=True, c=OCEAN, a=PP_ALIGN.CENTER)
    tb(s, 3.2, 3.18, 3.6, 0.4, "✓ COMPLETED", sz=13, b=True, c=OK, a=PP_ALIGN.CENTER)
    tb(s, 0.4, 4.15, 9.2, 0.5, "⭐ ⭐ ⭐ ⭐ ⭐ ⭐   6 颗 星 都 拿 到 啦!", sz=20, a=PP_ALIGN.CENTER)
    tb(s, 0.4, 4.78, 9.2, 0.3, "学 会 了 6 种 塑 料 · 6 个 替 代 · 我 会 认 6 词 · 我 会 写 3 词", sz=11, c=GRAY, a=PP_ALIGN.CENTER)

    page(share_close(prs, OCEAN,
        ["我 是 海 洋 小 卫 士, 我 可 以 ______。", "____ 不 好, 我 应 该 ____。"],
        "I'm an Ocean Guardian — I can ___ / I should ___.",
        "Day 4 — 家 庭 零 废 弃", "Family Zero Waste — 在 家 怎 样 零 废 弃?", "🏠"))

    for i, s in enumerate(prs.slides, 1):
        pn(s, i)
    prs.save("PPT/day3_water_plastic.pptx")
    print(f"Saved {n[0]} slides -> PPT/day3_water_plastic.pptx")


if __name__ == "__main__":
    build()
