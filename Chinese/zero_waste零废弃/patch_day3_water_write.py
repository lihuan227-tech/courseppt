#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Rewrite the three 我会写 slides (idx 26-28) of PPT/day3_water_plastic.pptx
into the day4_emergency format, IN PLACE (preserves the deck's embedded photos)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DECK = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PPT", "day3_water_plastic.pptx")
DEEP = RGBColor(0x0B,0x55,0x63); CREAM = RGBColor(0xFB,0xF7,0xEC); WARM = RGBColor(0xFF,0xF6,0xE6)
WHITE = RGBColor(0xFF,0xFF,0xFF); DARK = RGBColor(0x2C,0x2C,0x2C); GRAY = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB); FONT = 'KaiTi'

prs = Presentation(DECK)
W, H = prs.slide_width, prs.slide_height

def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT

def clear(s):
    spTree=s.shapes._spTree
    for el in list(spTree):
        if el.tag.split('}')[-1] in ('sp','pic','cxnSp','graphicFrame'):
            spTree.remove(el)

def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)

def build_write(s, w, py, en, tot, color=DEEP):
    clear(s); bg(s, CREAM)
    hd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.18),Inches(9.4),Inches(0.55))
    hd.fill.solid(); hd.fill.fore_color.rgb=color; hd.line.fill.background()
    tb(s,0.45,0.25,9.1,0.45,f"✍️ 我会写 · {w}  I Can Write",sz=16,b=True,c=WHITE)
    # LEFT big-char card
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(3.4))
    sh.fill.solid(); sh.fill.fore_color.rgb=WARM; sh.line.color.rgb=color; sh.line.width=Pt(2.5)
    csz={1:118,2:96,3:72}.get(len(w),72)
    tb(s,0.5,1.15,4.2,1.65,w,sz=csz,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.0,4.2,0.45,f"{py}  ·  {en}",sz=20,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.55,4.2,0.4,f"{tot} 笔 / {tot} strokes",sz=16,b=True,c=color,a=PP_ALIGN.CENTER)
    # RIGHT 3-step
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.6))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=WHITE; sh2.line.color.rgb=color; sh2.line.width=Pt(2)
    tb(s,5.15,1.1,4.4,0.4,"✏️ 3 步练习  3 Steps",sz=16,b=True,c=color)
    tb(s,5.15,1.55,4.4,0.35,"1️⃣ 看老师写  Watch teacher",sz=13,c=DARK)
    tb(s,5.15,1.90,4.4,0.35,"2️⃣ 用手指空中写  Air-write",sz=13,c=DARK)
    tb(s,5.15,2.25,4.4,0.35,"3️⃣ 在田字格写 3 次",sz=13,c=DARK)
    for i in range(4):
        x=5.0+i*1.15
        sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.85),Inches(1.05),Inches(1.05))
        sq.fill.solid(); sq.fill.fore_color.rgb=WHITE; sq.line.color.rgb=color; sq.line.width=Pt(1.5)
        ln1=s.shapes.add_connector(1,Inches(x),Inches(3.375),Inches(x+1.05),Inches(3.375)); ln1.line.color.rgb=LGRAY; ln1.line.width=Pt(0.5); ln1.line.dash_style=2
        ln2=s.shapes.add_connector(1,Inches(x+0.525),Inches(2.85),Inches(x+0.525),Inches(3.9)); ln2.line.color.rgb=LGRAY; ln2.line.width=Pt(0.5); ln2.line.dash_style=2
    tb(s,5.0,3.95,4.6,0.3,"在田字格里写 3 次 ↓",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.5),Inches(9.2),Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb=WARM; bar.line.color.rgb=color; bar.line.width=Pt(1.5)
    tb(s,0.55,4.55,4.4,0.24,"👩‍🏫 老师问 Teacher asks:",sz=10,b=True,c=color)
    tb(s,0.55,4.80,4.4,0.26,f"和我一起写「{w}」",sz=12,b=True,c=DARK)
    tb(s,5.05,4.55,4.5,0.24,"🧒 学生 Student does:",sz=10,b=True,c=color)
    tb(s,5.05,4.80,4.5,0.26,"看 → 空中写 → 田字格写 3 次",sz=12,b=True,c=DARK)

WRITE=[("水","shuǐ","water",4),("保护","bǎo hù","protect",16),("塑料","sù liào","plastic",23)]
for idx,data in zip([26,27,28],WRITE):
    build_write(prs.slides[idx], *data)

prs.save(DECK)
print("PATCHED 我会写 slides 26-28 in", DECK)
