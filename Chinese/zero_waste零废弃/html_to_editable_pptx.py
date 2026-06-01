#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
html_to_editable_pptx.py  —  Zero Waste · Day 1

Convert the self-contained HTML deck (day1_trash_version1.html) into a
FULLY EDITABLE 16:9 .pptx — native text boxes + colored rounded-rectangle
shapes (NO flat screenshots). Every slide's headings, cards, bins, steps,
goals, sentence-frames etc. stay selectable / editable in PowerPoint.

Pipeline:
  1. Headless-Chrome renders the deck so the 54 JS-generated slides exist as
     real DOM, then `--dump-dom` serializes the post-JS HTML.
  2. BeautifulSoup parses each `.slide`; a small recursive layout engine maps
     the component vocabulary (.zone/.goal/.step/.card/.framebar/...) onto
     native pptx shapes, resolving the CSS custom-properties + color-mix()
     to real colors.

Output: day1_trash_version1_editable.pptx  (NEW file — never overwrites).
Run:    python3 html_to_editable_pptx.py
"""
import os, re, subprocess, tempfile, shutil, html as _html
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SRC  = os.path.join(HERE, "day1_trash_version1.html")
OUT  = os.path.join(HERE, "day1_trash_version1_editable.pptx")
WORK = "/tmp/zw_edit"
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

# ---------------------------------------------------------------- colors
VARS = {
    'ink':'#0F1A3A','dark':'#2C2C2C','gray':'#88888c','lgray':'#BBBBBB',
    'green':'#2C5F2D','moss':'#6BA03D','teal':'#006970','leaf':'#C8E1A2',
    'cream':'#F6F8EC','warm':'#FFF3E0','imgbg':'#E8EEE0','gold':'#FFC107',
    'silver':'#B0B0B0','bronze':'#CD7F32','ok':'#2E7D32','alert':'#C8253E',
    'star':'#F5C242','b-recycle':'#1B6FBA','b-food':'#6BA03D',
    'b-harm':'#C8253E','b-other':'#88888c',
}
def _hex(h):
    h=h.strip().lstrip('#')
    if len(h)==3: h=''.join(c*2 for c in h)
    return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def _mix(a,b,pa):
    """color-mix(in srgb, a pa%, b)"""
    return RGBColor(*[round(x*pa+y*(1-pa)) for x,y in zip(a,b)])

WHITE=RGBColor(0xFF,0xFF,0xFF)
def lighten(c,amt=.82): return _mix(tuple(c),(255,255,255),amt)

def resolve_color(css, default='#2C2C2C'):
    """Resolve a CSS color token (var/hex/named/color-mix/gradient end) to RGBColor."""
    if not css: return _hex(default)
    s=css.strip()
    # gradient → take the LAST color token (saturated end)
    if 'gradient' in s:
        inner=s[s.find('(')+1:]
        # grab all var()/#hex/color-mix groups, use the last
        toks=re.findall(r'var\(--[\w-]+\)|#[0-9a-fA-F]{3,6}|color-mix\([^)]*\)|rgba?\([^)]*\)', inner)
        if toks: return resolve_color(toks[-1], default)
        return _hex(default)
    m=re.match(r'color-mix\(in srgb,\s*(.+?)\s+(\d+)%\s*,\s*(.+?)\)', s)
    if m:
        a=resolve_color(m.group(1),default); p=int(m.group(2))/100.0
        b=resolve_color(m.group(3),default)
        return _mix(tuple(a),tuple(b),p)
    m=re.match(r'var\(--([\w-]+)\)', s)
    if m: return _hex(VARS.get(m.group(1),default))
    if s.startswith('#'): return _hex(s)
    named={'white':'#FFFFFF','black':'#000000','#fff':'#FFFFFF'}
    if s.lower() in named: return _hex(named[s.lower()])
    if re.match(r'#?[0-9a-fA-F]{3,6}$',s): return _hex(s)
    return _hex(default)

def style_of(el):
    return el.get('style','') if el else ''
def css_prop(style, prop):
    m=re.search(rf'(?:^|;)\s*{prop}\s*:\s*([^;]+)', style)
    return m.group(1).strip() if m else None

# ---------------------------------------------------------------- geometry (inches)
SW, SH = 13.333, 7.5
def IN(v): return Emu(int(v*914400))

class Box:
    __slots__=('x','y','w','h')
    def __init__(s,x,y,w,h): s.x,s.y,s.w,s.h=x,y,w,h
    def pad(s,p): return Box(s.x+p,s.y+p,s.w-2*p,s.h-2*p)

# ---------------------------------------------------------------- text helpers
DISP="Noto Sans SC"   # display + body share a CJK-safe family for portability
BODY="Noto Sans SC"

def clean(t):
    return re.sub(r'\s+',' ', t).strip()

def text_lines(el):
    """Flatten element text into lines, breaking at <br> and block boundaries."""
    lines=['']
    def walk(n):
        for c in n.children:
            nm=getattr(c,'name',None)
            if nm is None:
                lines[-1]+=str(c)
            elif nm=='br':
                lines.append('')
            elif nm in ('div','p','section','ul','li','h4'):
                if lines[-1].strip(): lines.append('')
                walk(c)
                if lines[-1].strip(): lines.append('')
            else:
                walk(c)
    walk(el)
    return [clean(x) for x in lines if clean(x)]

def direct_text(el):
    """Text of element excluding nested block children, joining <br> as space."""
    parts=[]
    for c in el.children:
        if isinstance(c,NavigableString): parts.append(str(c))
        elif c.name=='br': parts.append(' ')
        elif c.name in ('span','b','small','i'): parts.append(c.get_text(' '))
    return clean(''.join(parts))

def add_text(slide, box, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, fill=None, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,italic,font)"""
    tb=slide.shapes.add_textbox(IN(box.x),IN(box.y),IN(box.w),IN(box.h))
    tf=tb.text_frame; tf.word_wrap=wrap
    tf.vertical_anchor=anchor
    for m in (tf.margin_left,):
        pass
    tf.margin_left=IN(0.05); tf.margin_right=IN(0.05)
    tf.margin_top=IN(0.02); tf.margin_bottom=IN(0.02)
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb=fill
        tb.line.fill.background()
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.alignment=align; p.line_spacing=line_spacing
        for (txt,size,color,bold,italic,font) in para:
            r=p.add_run(); r.text=txt
            r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
            r.font.name=font or BODY; r.font.color.rgb=color
            # set east-asian font too
            rPr=r._r.get_or_add_rPr()
            ea=rPr.find(qn('a:ea'))
            if ea is None:
                ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
            ea.set('typeface', font or BODY)
    return tb

def rrect(slide, box, fill, line=None, line_w=None, radius=0.10):
    sp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        IN(box.x),IN(box.y),IN(box.w),IN(box.h))
    try: sp.adjustments[0]=radius
    except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=line; sp.line.width=Pt(line_w or 1.5)
    sp.shadow.inherit=False
    sp.text_frame.word_wrap=True
    return sp

def R(text,size,color,bold=False,italic=False,font=None):
    return (text,size,color,bold,italic,font)

# colors as tuples for shorthand
C_INK=_hex(VARS['ink']); C_GREEN=_hex(VARS['green']); C_TEAL=_hex(VARS['teal'])
C_GRAY=_hex(VARS['gray']); C_MOSS=_hex(VARS['moss']); C_LEAF=_hex(VARS['leaf'])
C_CREAM=_hex(VARS['cream']); C_WARM=_hex(VARS['warm']); C_ALERT=_hex(VARS['alert'])
C_STAR=_hex(VARS['star']); C_ORANGE=_hex('#E65C00')

# ---------------------------------------------------------------- slide background
def bg(slide, color=C_CREAM):
    sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,IN(SW),IN(SH))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background()
    sp.shadow.inherit=False
    # dashed inner frame (echo .slide::before)
    fr=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,IN(.22),IN(.22),IN(SW-.44),IN(SH-.44))
    fr.fill.background(); fr.line.color.rgb=_mix(tuple(C_GREEN),(255,255,255),.35)
    fr.line.width=Pt(1.25); fr.shadow.inherit=False
    ln=fr.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    return sp

# move shape to back
def to_back(slide, sp):
    spTree=slide.shapes._spTree; spTree.remove(sp._element); spTree.insert(2,sp._element)

# ================================================================ component leaves
def has(el,*cls):
    c=set(el.get('class',[])); return any(k in c for k in cls)

def gettext(el,sel):
    n=el.select_one(sel); return clean(n.get_text(' ')) if n else ''

def emoji_and_rest(s):
    return s

def draw_zone(slide, el, box):
    col=resolve_color(css_prop(style_of(el),'background'), VARS['b-other'])
    rrect(slide, box, col, radius=0.10)
    ze=gettext(el,'.ze'); zn=gettext(el,'.zn'); zen=gettext(el,'.zen')
    rule=gettext(el,'.rule'); lis=[clean(li.get_text(' ')) for li in el.select('li')]
    paras=[]
    if ze: paras.append([R(ze,30,WHITE,True)])
    if zn: paras.append([R(zn,17,WHITE,True,font=DISP)])
    if zen: paras.append([R(zen,10,WHITE,False,True)])
    for li in lis: paras.append([R('• '+li,11,WHITE)])
    add_text(slide, box.pad(0.12), paras, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.TOP if lis else MSO_ANCHOR.MIDDLE)
    if rule:
        rh=0.5
        rb=Box(box.x+0.14, box.y+box.h-rh-0.12, box.w-0.28, rh)
        rrect(slide, rb, WHITE, radius=0.25)
        add_text(slide, rb, [[R(rule,10.5,C_INK,True)]], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

def draw_goal(slide, el, box):
    rrect(slide, box, WHITE, line=_hex('#eadfca'), line_w=1.5, radius=0.18)
    gn=gettext(el,'.gn'); gt=gettext(el,'.gt'); ge=gettext(el,'.ge'); gi=gettext(el,'.gi')
    bw=min(box.h-0.18, 0.62)
    bb=Box(box.x+0.14, box.y+(box.h-bw)/2, bw, bw)
    rrect(slide, bb, C_GREEN, radius=0.22)
    add_text(slide, bb, [[R(gn,22,WHITE,True,font=DISP)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx=box.x+0.14+bw+0.18
    tw=box.w-(0.14+bw+0.18)-1.0
    add_text(slide, Box(tx, box.y, tw, box.h),
             [[R(gt,15,C_INK,True)],[R(ge,10.5,C_GRAY,False,True)]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    if gi:
        add_text(slide, Box(box.x+box.w-1.0, box.y, 0.88, box.h),
                 [[R(gi,20,C_INK)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def draw_step(slide, el, box):
    rrect(slide, box, WHITE, line=_hex('#eadfca'), line_w=1.5, radius=0.16)
    n=gettext(el,'.n'); e=gettext(el,'.e'); t=gettext(el,'.t'); d=gettext(el,'.d')
    add_text(slide, box.pad(0.10),
             [[R(n,20,C_MOSS,True,font=DISP)],[R(e,30,C_INK)],
              [R(t,15,C_INK,True)],[R(d,10,C_GRAY,False,True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def draw_it(slide, el, box):
    bc=resolve_color(css_prop(style_of(el),'border-color'), VARS['leaf'])
    rrect(slide, box, WHITE, line=bc, line_w=2, radius=0.16)
    e=gettext(el,'.e'); l=gettext(el,'.l'); p=gettext(el,'.p')
    add_text(slide, box.pad(0.08),
             [[R(e,30,C_INK)],[R(l,14,C_INK,True)],[R(p,10,C_GRAY,False,True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def draw_qcard(slide, el, box):
    rrect(slide, box, WHITE, line=C_LEAF, line_w=2, radius=0.18)
    q=gettext(el,'.q'); qen=gettext(el,'.qen')
    # .q may itself contain .qen — separate
    qnode=el.select_one('.q')
    if qnode:
        qen2=qnode.select_one('.qen'); 
        if qen2: qen=clean(qen2.get_text(' ')); 
        q=clean(''.join(s for s in qnode.find_all(string=True, recursive=True)
                        if not (qen2 and qen2 in getattr(s,'parents',[]))))
        q=clean(qnode.get_text(' '))
        if qen and qen in q: q=clean(q.replace(qen,''))
    paras=[[R('❓ '+q,15,C_INK,True)]]
    if qen: paras.append([R(qen,11,C_GRAY,False,True)])
    add_text(slide, box.pad(0.16), paras, anchor=MSO_ANCHOR.MIDDLE)

def draw_framebar(slide, el, box):
    col=resolve_color(css_prop(style_of(el),'background'), VARS['green'])
    rrect(slide, box, col, radius=0.16)
    tag=gettext(el,'.tag'); en=gettext(el,'.en')
    frnode=el.select_one('.frame')
    fr=''
    if frnode:
        ennode=frnode.select_one('.en')
        if ennode and not en: en=clean(ennode.get_text(' '))
        if ennode: ennode.extract()
        fr=clean(frnode.get_text(' '))
    th=box.h-0.2
    if tag:
        tb=Box(box.x+0.16, box.y+(box.h-0.5)/2, 1.5, 0.5)
        rrect(slide, tb, _mix(tuple(col),(255,255,255),.78), radius=0.5)
        add_text(slide, tb, [[R(tag,11,WHITE,True,font=DISP)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx=box.x+0.16+1.5+0.18
    else:
        tx=box.x+0.2
    paras=[[R(fr,16,WHITE,True)]]
    if en: paras.append([R(en,11,WHITE,False,True)])
    add_text(slide, Box(tx, box.y, box.x+box.w-tx-0.16, box.h), paras, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

def draw_safety(slide, el, box):
    rrect(slide, box, C_ALERT, radius=0.14)
    ic=gettext(el,'.ic'); en=gettext(el,'.en')
    ennode=el.select_one('.en')
    if ennode: ennode.extract()
    txt=clean(el.get_text(' '))
    if ic and txt.startswith(ic): txt=clean(txt[len(ic):])
    paras=[[R(('🛑 ' if not ic else ic+' ')+txt,16,WHITE,True)]]
    if en: paras.append([R(en,11,WHITE,False,True)])
    add_text(slide, box.pad(0.16), paras, anchor=MSO_ANCHOR.MIDDLE)

def draw_timer(slide, el, box):
    rrect(slide, box, C_WARM, line=_hex('#f3d8a8'), line_w=2, radius=0.18)
    big=gettext(el,'.big'); lab=gettext(el,'.lab')
    add_text(slide, box.pad(0.08), [[R(big,30,C_ORANGE,True,font=DISP)],[R(lab,12,C_ORANGE,True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def draw_medal(slide, el, box):
    col=resolve_color(css_prop(style_of(el),'background'), VARS['gold'])
    tc=resolve_color(css_prop(style_of(el),'color'), '#5a4500')
    rrect(slide, box, col, radius=0.18)
    m=gettext(el,'.m'); mt=gettext(el,'.mt'); md=gettext(el,'.md')
    add_text(slide, box.pad(0.12), [[R(m,40,tc)],[R(mt,16,tc,True,font=DISP)],[R(md,11,tc)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def draw_minibin(slide, el, box):
    correct=has(el,'correct')
    if correct:
        col=resolve_color(css_prop(style_of(el),'background'), VARS['green'])
        rrect(slide, box, col, radius=0.14)
        tc=WHITE
    else:
        rrect(slide, box, WHITE, line=_hex('#e3dcc8'), line_w=2, radius=0.14)
        tc=_hex('#c9c4b4')
    mbe=gettext(el,'.mbe'); mbn=gettext(el,'.mbn')
    tick='  ✓' if correct else ''
    add_text(slide, box.pad(0.08), [[R(mbe,24,tc)],[R(mbn+tick,15,tc,True,font=DISP)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def draw_bigitem(slide, el, box):
    sp=rrect(slide, box, WHITE, line=_hex(VARS['lgray']), line_w=2.5, radius=0.18)
    ln=sp.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    be=gettext(el,'.be'); bl=gettext(el,'.bl'); ben=gettext(el,'.ben')
    add_text(slide, box.pad(0.12), [[R(be,66,C_INK)],[R(bl,24,C_INK,True,font=DISP)],[R(ben,12,C_GRAY,False,True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def draw_photo(slide, el, box):
    sp=rrect(slide, box, _hex(VARS['imgbg']), line=_hex(VARS['lgray']), line_w=2.5, radius=0.18)
    ln=sp.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    txt=clean(el.get_text(' '))
    small=el.select_one('small'); sm=''
    if small: sm=clean(small.get_text(' ')); txt=clean(txt.replace(sm,''))
    txt=txt.replace('📷','').strip()
    paras=[[R('📷',34,C_GRAY)],[R(txt,14,C_GRAY,True)]]
    if sm: paras.append([R(sm,10,_hex(VARS['lgray']))])
    add_text(slide, box.pad(0.12), paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

def draw_vidwrap(slide, el, box):
    rrect(slide, box, _hex('#0b1410'), radius=0.14)
    vtag=gettext(el,'.vtag')
    if vtag:
        tb=Box(box.x+0.14, box.y+0.14, min(box.w-0.28,4.2), 0.42)
        rrect(slide, tb, _hex('#0F1A3A'), radius=0.5)
        add_text(slide, tb, [[R(vtag,11,WHITE,True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, box.pad(0.2),
             [[R('▶',40,_hex('#a7e0a8'))],[R('点击播放视频',18,_hex('#a7e0a8'),True)],
              [R('Click to play (live deck)',11,_hex('#7fb98a'))]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

def draw_chip(slide, el, box):
    style=style_of(el)
    bgc=css_prop(style,'background')
    fill=resolve_color(bgc, 'white') if bgc else WHITE
    tcol=resolve_color(css_prop(style,'color'),'#2C5F2D') if css_prop(style,'color') else C_GREEN
    line=None
    if not bgc or 'var(--warm)' in (bgc or '') or fill==WHITE:
        line=C_LEAF
    rrect(slide, box, fill, line=line, line_w=2, radius=0.5)
    add_text(slide, box.pad(0.08), [[R(clean(el.get_text(' ')),13,tcol,True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def draw_bigbtn(slide, el, box):
    rrect(slide, box, C_ORANGE, radius=0.2)
    add_text(slide, box.pad(0.1), [[R('▶  '+clean(el.get_text(' ')),18,WHITE,True,font=DISP)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def draw_tian(slide, el, box):
    side=min(box.w, box.h)
    bx=Box(box.x+(box.w-side)/2, box.y+(box.h-side)/2, side, side)
    sp=rrect(slide, bx, WHITE, line=C_ALERT, line_w=2.5, radius=0.04)
    gh=gettext(el,'.gh')
    if gh:
        add_text(slide, bx, [[R(gh,int(side*54),_mix(tuple(C_GREEN),(255,255,255),.20))]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def draw_card(slide, el, box, fill=WHITE):
    style=style_of(el)
    bgc=css_prop(style,'background')
    f=resolve_color(bgc,'#FFFFFF') if (bgc and 'gradient' not in bgc and 'var' in bgc) else WHITE
    if bgc and 'gradient' in bgc: f=resolve_color(bgc,'#FFFFFF')
    bl=css_prop(style,'border'); blcol=None
    border_left=css_prop(style,'border-left')
    rr=rrect(slide, box, f if f else WHITE, line=_hex('#eadfca'), line_w=1.5, radius=0.16)
    # accent left border
    if border_left:
        m=re.search(r'(#[0-9a-fA-F]{3,6}|var\(--[\w-]+\))', border_left)
        if m:
            acc=resolve_color(m.group(1))
            ab=Box(box.x, box.y+0.06, 0.10, box.h-0.12)
            rrect(slide, ab, acc, radius=0.4)
    # content: heading(s) + text
    inner=box.pad(0.18)
    paras=[]
    bt=el.select_one('.blockt')
    btcol=C_GREEN
    if bt:
        bs=style_of(bt)
        if css_prop(bs,'color'): btcol=resolve_color(css_prop(bs,'color'),VARS['green'])
        paras.append([R(clean(bt.get_text(' ')),16,btcol,True,font=DISP)])
        bt.extract()
    for line in text_lines(el):
        paras.append([R(line,14,C_INK)])
    if not paras:
        rest=clean(el.get_text(' '))
        paras=[[R(rest,14,C_INK)]]
    add_text(slide, inner, paras, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

# ================================================================ layout engine
LEAF = {
    'zone':draw_zone,'goal':draw_goal,'step':draw_step,'it':draw_it,
    'qcard':draw_qcard,'framebar':draw_framebar,'safety':draw_safety,
    'timer':draw_timer,'medal':draw_medal,'minibin':draw_minibin,
    'bigitem':draw_bigitem,'photo':draw_photo,'vidwrap':draw_vidwrap,
    'chip':draw_chip,'bigbtn':draw_bigbtn,'tian':draw_tian,'card':draw_card,
}
def leaf_kind(el):
    cs=el.get('class',[])
    for k in ('zone','goal','step','it','qcard','framebar','safety','timer',
              'medal','minibin','bigitem','photo','vidwrap','tian','bigbtn','chip','card'):
        if k in cs: return k
    return None

INLINE={'br','script','style','span','b','small','i','a','em','strong'}
def child_tags(el):
    return [c for c in el.children if getattr(c,'name',None) and c.name not in INLINE]

def is_row(el):
    cs=el.get('class',[]); st=style_of(el)
    if 'row' in cs or 'steps' in cs: return True
    if 'display:flex' in st and 'column' not in st: return True
    return False
def is_col(el):
    st=style_of(el)
    return ('display:flex' in st and 'column' in st)
def is_grid(el):
    return 'display:grid' in style_of(el) or 'minibins' in el.get('class',[]) or 'itemstrip' in el.get('class',[]) or 'goals' in el.get('class',[])

def flex_weight(el):
    st=style_of(el); cs=el.get('class',[])
    if 'arrow' in cs: return ('fixed',0.45)
    m=re.search(r'flex\s*:\s*([0-9.]+)', st)
    if m: return ('flex',float(m.group(1)))
    if re.search(r'flex\s*:\s*none', st): return ('fixed',None)
    return ('flex',1.0)

GAP=0.16

def render_block(slide, el, box, depth=0):
    """Recursively place an element tree into box."""
    if box.w<=0.2 or box.h<=0.15: return
    cs=el.get('class',[])
    k=leaf_kind(el)
    # treat as leaf only if it has no block children worth recursing (cards recurse internally)
    if k and k!='card':
        LEAF[k](slide, el, box); return
    if k=='card':
        # card with nested row/grid → render container inside; else text card
        kids=child_tags(el)
        if any(is_row(c) or is_grid(c) or leaf_kind(c) in ('zone','step','goal') for c in kids):
            rrect(slide, box, WHITE, line=_hex('#eadfca'), line_w=1.5, radius=0.16)
            inner=box.pad(0.16)
            stack_vertical(slide, kids, inner, depth+1)
        else:
            draw_card(slide, el, box)
        return
    # itemstrip → wrap
    if 'itemstrip' in cs:
        items=[c for c in child_tags(el) if leaf_kind(c)=='it' or 'it' in c.get('class',[])]
        wrap_items(slide, items, box, depth); return
    if 'goals' in cs:
        stack_vertical(slide, [c for c in child_tags(el) if 'goal' in c.get('class',[])], box, depth, gap=0.14); return
    if 'minibins' in cs:
        grid(slide, [c for c in child_tags(el) if 'minibin' in c.get('class',[])], box, cols=2, depth=depth); return
    if is_grid(el):
        cols=2
        m=re.search(r'grid-template-columns:\s*repeat\((\d+)', style_of(el))
        if m: cols=int(m.group(1))
        elif 'grid-template-columns:1fr 1fr' in style_of(el).replace(' ',''): cols=2
        grid(slide, child_tags(el), box, cols=cols, depth=depth); return
    if is_row(el):
        row_layout(slide, child_tags(el), box, depth); return
    if is_col(el):
        stack_vertical(slide, child_tags(el), box, depth); return
    # generic wrapper: if single/multiple block kids, stack; else text
    kids=child_tags(el)
    if kids:
        stack_vertical(slide, kids, box, depth)
    else:
        txt=clean(el.get_text(' '))
        if txt:
            add_text(slide, box, [[R(txt,15,C_INK)]], anchor=MSO_ANCHOR.MIDDLE)

def row_layout(slide, kids, box, depth):
    kids=[c for c in kids if clean(c.get_text(' ')) or leaf_kind(c) or child_tags(c)]
    if not kids: return
    weights=[flex_weight(c) for c in kids]
    fixed=sum(w for t,w in weights if t=='fixed' and w)
    nfix=sum(1 for t,w in weights if t=='fixed' and not w)
    fixed += nfix*1.4
    flexsum=sum(w for t,w in weights if t=='flex') or 1
    gaps=GAP*(len(kids)-1)
    avail=box.w-gaps-fixed
    x=box.x
    for c,(t,w) in zip(kids,weights):
        if t=='fixed':
            cw=w if w else 1.4
        else:
            cw=avail*(w/flexsum)
        if 'arrow' in c.get('class',[]):
            add_text(slide, Box(x,box.y,cw,box.h), [[R('→',26,C_MOSS,True)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            render_block(slide, c, Box(x,box.y,cw,box.h), depth+1)
        x+=cw+GAP

def stack_vertical(slide, kids, box, depth, gap=GAP):
    kids=[c for c in kids if clean(c.get_text(' ')) or leaf_kind(c) or child_tags(c)]
    if not kids: return
    # assign heights: bars/chips/timers fixed, others share
    def fixed_h(c):
        cs=c.get('class',[])
        if 'framebar' in cs: return 0.95
        if 'safety' in cs: return 0.85
        if 'timer' in cs: return 1.15
        if 'chip' in cs: return 0.55
        if leaf_kind(c)=='qcard': return None
        return None
    fixed=[fixed_h(c) for c in kids]
    # header-ish small text node?
    fsum=sum(h for h in fixed if h)
    flexn=[1 for h in fixed if h is None]
    # weight by flex style
    wts=[]
    for c,h in zip(kids,fixed):
        if h is not None: wts.append(None); continue
        t,w=flex_weight(c); 
        wts.append(w if t=='flex' else 1.0)
    fsum_gaps=gap*(len(kids)-1)
    avail=box.h-fsum-fsum_gaps
    wsum=sum(w for w in wts if w) or 1
    y=box.y
    for c,h,w in zip(kids,fixed,wts):
        ch=h if h is not None else avail*(w/wsum)
        render_block(slide, c, Box(box.x,y,box.w,ch), depth+1)
        y+=ch+gap

def wrap_items(slide, items, box, depth):
    if not items: return
    n=len(items)
    per_row=min(n, 5 if n<=5 else (n+1)//2)
    rows=(n+per_row-1)//per_row
    iw=(box.w-GAP*(per_row-1))/per_row
    ih=min((box.h-GAP*(rows-1))/rows, 1.7)
    total_h=ih*rows+GAP*(rows-1)
    y0=box.y+(box.h-total_h)/2
    for i,it in enumerate(items):
        r=i//per_row; cc=i%per_row
        in_row = per_row if r<rows-1 else (n-per_row*(rows-1))
        rw=(box.w-GAP*(in_row-1))/in_row if in_row>0 else iw
        rw=min(rw, 2.1)
        row_w=rw*in_row+GAP*(in_row-1)
        x0=box.x+(box.w-row_w)/2
        x=x0+cc*(rw+GAP)
        draw_it(slide, it, Box(x, y0+r*(ih+GAP), rw, ih))

def grid(slide, kids, box, cols, depth):
    kids=[c for c in kids if clean(c.get_text(' ')) or leaf_kind(c) or child_tags(c)]
    if not kids: return
    n=len(kids); rows=(n+cols-1)//cols
    cw=(box.w-GAP*(cols-1))/cols
    ch=(box.h-GAP*(rows-1))/rows
    for i,c in enumerate(kids):
        r=i//cols; cc=i%cols
        render_block(slide, c, Box(box.x+cc*(cw+GAP), box.y+r*(ch+GAP), cw, ch), depth+1)

# ================================================================ header / special slides
def render_header(slide, sec):
    """Render .kicker + .s-head wrapper; return content-top y."""
    wrap=None
    for d in sec.select('.anim'):
        if d.select_one('.kicker') or d.select_one('.s-head'): wrap=d; break
    if not wrap:
        sh=sec.select_one('.s-head')
        if not sh: return 1.5
    kicker=sec.select_one('.kicker'); shead=sec.select_one('.s-head')
    x=0.62; y=0.5; w=SW-1.24
    if kicker:
        add_text(slide, Box(x,y,w,0.34), [[R('— '+clean(kicker.get_text(' ')),13,C_TEAL,False,True,DISP)]])
        y+=0.36
    if shead:
        en=shead.select_one('.en'); ent=''
        if en: ent=clean(en.get_text(' ')); en.extract()
        ht=clean(shead.get_text(' '))
        add_text(slide, Box(x,y,w,0.66), [[R(ht,27,C_GREEN,True,font=DISP)]])
        y+=0.62
        if ent:
            add_text(slide, Box(x,y,w,0.3), [[R(ent,12,C_GRAY,False,True)]])
            y+=0.32
    return y+0.14

def render_cover(slide, sec):
    bg(slide)
    add_text(slide, Box(0,1.0,SW,0.9), [[R(clean(sec.select_one('.float-row').get_text(' ')),46,C_INK)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    band=sec.select_one('.cover-band')
    bw=5.6
    bb=Box((SW-bw)/2,2.05,bw,0.62)
    rrect(slide, bb, C_GREEN, radius=0.5)
    add_text(slide, bb, [[R(clean(band.get_text(' ')),19,WHITE,True,font=DISP)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    title=sec.select_one('.cover-title')
    add_text(slide, Box(0,2.85,SW,1.5), [[R(clean(title.get_text(' ')),66,C_GREEN,True,font=DISP)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    q=sec.select_one('.cover-q')
    en=q.select_one('.en'); ent=clean(en.get_text(' ')) if en else ''
    if en: en.extract()
    add_text(slide, Box(0,4.5,SW,0.9),
             [[R(clean(q.get_text(' ')),20,C_TEAL,True)],[R(ent,13,C_GRAY,False,True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    mt=sec.select_one('.monster-teaser')
    if mt:
        add_text(slide, Box(0,5.6,SW,0.6), [[R(clean(mt.get_text(' ')),15,C_ALERT,True)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def render_divider(slide, sec):
    col=resolve_color(css_prop(style_of(sec),'background'), VARS['green'])
    sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,IN(SW),IN(SH))
    sp.fill.solid(); sp.fill.fore_color.rgb=col; sp.line.fill.background(); sp.shadow.inherit=False
    de=sec.select_one('.de'); dt=sec.select_one('.dt'); ds=sec.select_one('.ds')
    paras=[]
    if de: paras.append([R(clean(de.get_text(' ')),54,WHITE)])
    if dt: paras.append([R(clean(dt.get_text(' ')),40,WHITE,True,font=DISP)])
    if ds: paras.append([R(clean(ds.get_text(' ')),17,WHITE,False)])
    extra=sec.select_one('.dt ~ div.anim') or None
    add_text(slide, Box(0.8,1.0,SW-1.6,4.6), paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    dtime=sec.select_one('.dtime')
    if dtime:
        pills=[clean(t.get_text(' ')) for t in dtime.select('.t')]
        pw=2.4; gap=0.3; tot=pw*len(pills)+gap*(len(pills)-1)
        x=(SW-tot)/2; yb=5.3
        for p in pills:
            b=Box(x,yb,pw,0.62); rr=rrect(slide,b,_mix(tuple(col),(255,255,255),.78),radius=0.3)
            add_text(slide,b,[[R(p,17,WHITE,True,font=DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            x+=pw+gap
    # trailing emoji row (group photo slide)
    for d in sec.select('.anim'):
        t=clean(d.get_text(' '))
        if d is not de and d is not dt and d is not ds and t and 'dtime' not in d.get('class',[]):
            add_text(slide, Box(0,5.4,SW,0.8),[[R(t,34,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            break

# ================================================================ main
def render_content(slide, sec):
    bg(slide)
    top=render_header(slide, sec)
    # body = top-level children after the header wrapper
    header_wrap=None
    for d in sec.find_all('div', recursive=False):
        if d.select_one('.kicker') or d.select_one('.s-head'): header_wrap=d; break
    body=[c for c in child_tags(sec) if c is not header_wrap and not (c.select_one('.kicker') or c.select_one('.s-head'))]
    box=Box(0.62, top, SW-1.24, SH-top-0.45)
    if len(body)==1:
        render_block(slide, body[0], box, 0)
    else:
        stack_vertical(slide, body, box, 0)

def main():
    os.makedirs(WORK, exist_ok=True)
    html=open(SRC,encoding='utf-8').read()
    render=re.sub(r'src="\$\{vsrc\([^)]*\)\}"','',html)
    rpath=os.path.join(WORK,'render.html'); open(rpath,'w',encoding='utf-8').write(render)
    dom=os.path.join(WORK,'dom.html')
    def have_dom(p):
        return os.path.exists(p) and os.path.getsize(p)>40000 and 'class="slide' in open(p,encoding='utf-8').read()
    if not have_dom(dom):
        prof=tempfile.mkdtemp(prefix='zwdom_')
        try:
            with open(dom,'w',encoding='utf-8') as f:
                proc=subprocess.Popen([CHROME,'--headless=new','--disable-gpu','--no-first-run',
                    '--no-default-browser-check',f'--user-data-dir={prof}','--dump-dom',
                    f'file://{rpath}','--virtual-time-budget=2500'],
                    stdout=f, stderr=subprocess.DEVNULL)
                try: proc.wait(timeout=90)
                except subprocess.TimeoutExpired: proc.kill()
        finally:
            shutil.rmtree(prof, ignore_errors=True)
        if not have_dom(dom):
            raise SystemExit('DOM render failed — dump Chrome --dump-dom into '+dom)
    soup=BeautifulSoup(open(dom,encoding='utf-8').read(),'lxml')
    slides=soup.select('#stage > .slide')
    print(f'parsed {len(slides)} slides')

    prs=Presentation(); prs.slide_width=IN(SW); prs.slide_height=IN(SH)
    blank=prs.slide_layouts[6]
    for i,sec in enumerate(slides):
        s=prs.slides.add_slide(blank)
        cs=sec.get('class',[])
        try:
            if sec.get('id')=='cover': render_cover(s, sec)
            elif 'divider' in cs: render_divider(s, sec)
            else: render_content(s, sec)
        except Exception as e:
            bg(s)
            add_text(s, Box(0.7,3,SW-1.4,1.5), [[R(f'[slide {i+1} fallback] '+clean(sec.get_text(' '))[:300],12,C_INK)]],
                     anchor=MSO_ANCHOR.MIDDLE)
            print(f'  !! slide {i+1}: {type(e).__name__}: {e}')
        if (i+1)%10==0: print(f'  ...{i+1}/{len(slides)}')
    prs.save(OUT)
    print('saved', OUT, len(prs.slides._sldIdLst),'slides')

if __name__=='__main__':
    main()
