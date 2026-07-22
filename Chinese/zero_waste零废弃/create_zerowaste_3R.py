#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零废弃与 3R — Zero Waste & the 3R (Reduce · Reuse · Recycle)
Story-driven K-5 immersive lesson, built on the 野外生存与探险 Day-3 structure.

Session 1 (~46 min):
  绘本导入《如果地球被我们吃掉》 → 什么是 Zero Waste / 3R → 3R 分类游戏
  → Zero Waste 大侦探 → 学校 Zero Waste 讨论 → 变废为宝视频 + Project 过渡
Session 2 (~40 min):
  复习游戏 → 指对字 → 我会认 (环保/垃圾/回收/减少/地球) → 我会写 (环 · 保)
Session 3 (~60 min):
  变废为宝 · 环保艺术展 — 3 leveled projects (拼贴画 / 瓶子变变变 / 环保海报·发明)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Zero-Waste palette (leafy green primary, 3 distinct R accents) ---
ECO     = RGBColor(0x2E,0x7D,0x32)   # primary green (headers)
LEAF    = RGBColor(0x66,0xBB,0x6A)   # light green
DEEP    = RGBColor(0x1B,0x5E,0x20)   # deep green (cover / dividers)
REDUCE  = RGBColor(0x1E,0x88,0xE5)   # blue  — Reduce 减少
REUSE   = RGBColor(0xFB,0x8C,0x00)   # orange — Reuse 重复使用
RECYCLE = RGBColor(0x43,0xA0,0x47)   # green — Recycle 回收
SUN     = RGBColor(0xF5,0xC2,0x42)
BROWN   = RGBColor(0x6B,0x44,0x23)
SKY     = RGBColor(0x4A,0x90,0xD9)
ALERT   = RGBColor(0xD0,0x4A,0x3C)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xF1,0xF8,0xE9)   # pale green wash (cards)
CREAM   = RGBColor(0xFD,0xFB,0xF3)
IMGBG   = RGBColor(0xE8,0xE8,0xE8)
GOLD    = RGBColor(0xF9,0xA8,0x25)
OK      = RGBColor(0x38,0x8E,0x3C)

BASE = "/Users/Huan/0 projects/summercourse/Chinese/zero_waste零废弃"
STORY_IMG = os.path.join(BASE, "pics", "zw_earth_eaten.png")

# === Helpers (identical scaffolding to 野外生存与探险 Day 3) ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def img(s,l,t,w,h,path,fallback="📷"):
    if os.path.exists(path):
        s.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(h))
    else:
        ib(s,l,t,w,h,fallback)
def hb(s,txt,c=ECO,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.8,9,0.8,sub,sz=20,c=WHITE,a=PP_ALIGN.CENTER);return s
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)
def teacher_student_bar(s,t,teacher_q,student_action,color=ECO):
    """Bottom prompt bar — teacher question on left, student action on right."""
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=color;sf.line.width=Pt(2)
    tb(s,0.45,t+0.04,4.5,0.25,"👩‍🏫 老师问 Teacher asks:",sz=10,b=True,c=color)
    tb(s,0.45,t+0.27,4.5,0.28,teacher_q,sz=12,b=True,c=DARK)
    tb(s,5.0,t+0.04,4.6,0.25,"🧒 学生 Student does:",sz=10,b=True,c=REUSE)
    tb(s,5.0,t+0.27,4.6,0.28,student_action,sz=12,b=True,c=DARK)

def word_card_read(w,py,en,img_emoji,sent,color=ECO):
    """Image-driven recognition — emoji LEFT, character+pinyin RIGHT, sentence bottom."""
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(2.6))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,0.5,1.3,4.2,2.0,img_emoji,sz=130,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(2.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=color;sh2.line.width=Pt(2.5)
    char_sz = 110 if len(w)==1 else (66 if len(w)==2 else 54)
    tb(s,5.1,1.15,4.4,1.5,w,sz=char_sz,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,5.1,2.70,4.4,0.4,f"{py}  ·  {en}",sz=18,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,5.1,3.15,4.4,0.4,"👉 跟我读 3 次！",sz=14,c=color,a=PP_ALIGN.CENTER)
    sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.85),Inches(9.2),Inches(1.3))
    sh3.fill.solid();sh3.fill.fore_color.rgb=WHITE;sh3.line.color.rgb=color;sh3.line.width=Pt(2)
    tb(s,0.6,3.95,2.0,0.4,"💬 例句 Example",sz=14,b=True,c=color)
    tb(s,0.6,4.35,8.8,0.6,sent,sz=22,b=True,c=DARK)
    return s

def word_card_write(w,py,en,strokes_count,strokes_hint,color=ECO):
    """Watch → air-write → write in 田字格."""
    s=ns();bg(s,CREAM);hb(s,"✍️ 我会写  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(4.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,0.5,1.1,4.2,2.4,w,sz=160,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.55,4.2,0.5,f"{py}  ·  {en}",sz=22,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.10,4.2,0.4,f"{strokes_count} 笔 / {strokes_count} strokes",sz=16,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.50,4.2,0.4,strokes_hint,sz=11,c=DARK,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,5.15,1.1,4.4,0.4,"3 步练习  3 Steps",sz=16,b=True,c=color)
    tb(s,5.15,1.5,4.4,0.4,"1️⃣ 看老师写  Watch teacher",sz=13,c=DARK)
    tb(s,5.15,1.85,4.4,0.4,"2️⃣ 用手指空中写  Air-write",sz=13,c=DARK)
    tb(s,5.15,2.20,4.4,0.4,"3️⃣ 在田字格写 3 次  Write 3 times",sz=13,c=DARK)
    for i in range(4):
        x=5.0+i*1.15
        sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.85),Inches(1.05),Inches(1.05))
        sq.fill.solid();sq.fill.fore_color.rgb=WHITE;sq.line.color.rgb=color;sq.line.width=Pt(1.5)
        ln1=s.shapes.add_connector(1,Inches(x),Inches(2.85+0.525),Inches(x+1.05),Inches(2.85+0.525))
        ln1.line.color.rgb=LGRAY;ln1.line.width=Pt(0.5);ln1.line.dash_style=2
        ln2=s.shapes.add_connector(1,Inches(x+0.525),Inches(2.85),Inches(x+0.525),Inches(2.85+1.05))
        ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.5);ln2.line.dash_style=2
    tb(s,5.0,4.0,4.6,0.3,"在田字格里写 3 次 ↓  Write here 3 times",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    teacher_student_bar(s,4.45,f"和我一起写「{w}」","看 → 空中写 → 田字格写 3 次")
    return s

def r_example_slide(letter,en_r,cn_r,py,meaning,color,examples,keep_line):
    """One of the 3R — big meaning + 4 example cards."""
    s=ns();bg(s,CREAM);hb(s,f"{letter}  {en_r} · {cn_r}  ({py})",color)
    tb(s,0.4,0.85,9.2,0.45,meaning,sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
    for i,(em,cap) in enumerate(examples):
        x=0.4+i*2.35
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.55),Inches(2.20),Inches(2.35))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
        tb(s,x+0.05,1.75,2.1,1.0,em,sz=60,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,2.95,2.1,0.85,cap,sz=13,b=True,c=color,a=PP_ALIGN.CENTER)
    kb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.05),Inches(9.2),Inches(0.5))
    kb.fill.solid();kb.fill.fore_color.rgb=color;kb.line.fill.background()
    tb(s,0.4,4.12,9.2,0.4,keep_line,sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    return s

# ========================================================================
#                              SLIDES
# ========================================================================
n=0

# 1. COVER
s=ns();bg(s,DEEP)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.4,8,0.5,"ZERO WASTE",sz=18,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"♻️ 零废弃与 3R",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"Zero Waste  ·  Reduce · Reuse · Recycle",sz=20,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"🌍 环保小卫士任务  Earth Guardian Mission",sz=24,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"读绘本 · 认 3R · 分类 · 当侦探 · 变废为宝",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"Read · Sort · Investigate · Discuss · Create",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"零废弃 · Zero Waste Unit",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「环保小卫士们! 我们的地球生病了 — 今天我们要一起帮它!」\n• 不先讲道理 — 先读绘本, 让学生自己发现问题。")

# 2. SESSION 1 GOALS
s=ns();bg(s,CREAM);hb(s,"🎯 Session 1 学习目标  Today's Goals",ECO)
tb(s,0.4,0.85,9.2,0.4,"上完 Session 1, 你可以…",sz=18,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"After Session 1, you can…",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
goals=[
    ("1️⃣","知道地球需要我们保护","Know the Earth needs us (not money!)",OK),
    ("2️⃣","知道什么是「零废弃」","Know what Zero Waste means",ECO),
    ("3️⃣","认识 3R","Reduce 减少 · Reuse 重复用 · Recycle 回收",RECYCLE),
    ("4️⃣","把垃圾分到 3R","Sort trash into the 3R",REUSE),
    ("5️⃣","找出不环保, 说更好的做法","Spot the problem, say a better way",REDUCE),
]
for i,(num,cn,en,cl) in enumerate(goals):
    y=1.70+i*0.65
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.58))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.55,y+0.10,0.7,0.4,num,sz=22,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.05,3.9,0.3,cn,sz=17,b=True,c=cl)
    tb(s,5.30,y+0.10,4.2,0.3,en,sz=12,c=DARK)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「上 Session 1 我们要学会 5 件事 — 一起读!」\n• 跟读 5 个目标, 完成 1 个打 1 个 ✓。\n• 关键词: 保护地球 / 零废弃 / 3R / 分类 / 大侦探。")

# 3. SESSION 1 DIVIDER
s=div("Session 1  上午 46 min","📖 绘本 · 3R · 分类 · 大侦探 · 变废为宝",ECO,"🌍");n+=1;pn(s,n)

# 4. STORY — picture book hook (8 min · part 1)
s=ns();bg(s,CREAM);hb(s,"📖 绘本时间  Story Time",ECO)
img(s,0.3,0.95,5.0,3.6,STORY_IMG,"📷 绘本《如果地球被我们吃掉》插图 / picture-book illustration")
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.5),Inches(0.95),Inches(4.2),Inches(3.6))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=ECO;sh.line.width=Pt(2.5)
tb(s,5.7,1.15,3.8,0.5,"《如果地球被我们吃掉》",sz=18,b=True,c=ECO)
tb(s,5.7,1.75,3.8,0.4,"If We Ate Up the Whole Earth",sz=12,c=GRAY)
tf=tb(s,5.7,2.25,3.8,0.5,"🌊 水  Water",sz=17,b=True,c=REDUCE)
ap(tf,"🌬️ 空气  Air",sz=17,b=True,c=SKY)
ap(tf,"🌳 树木  Trees",sz=17,b=True,c=ECO)
ap(tf,"🐾 动物  Animals",sz=17,b=True,c=BROWN)
tb(s,0.4,4.65,9.2,0.4,"我们一边看图, 一边想: 图里有什么? 它们重要吗?",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"绘本导入 (8 分钟 · 第 1 部分):\n• 慢慢看图讲绘本《如果地球被我们吃掉》。\n• 边看边指: 水 / 空气 / 树木 / 动物 — 它们给我们什么?\n• 先不给结论 — 下一页提出思考问题, 让学生自己想。")

# 5. STORY — think question (think-pair-share, no answer yet)
s=ns();bg(s,CREAM);hb(s,"🤔 想一想  Big Question",SUN)
tb(s,0.4,0.90,9.2,0.55,"如果水、空气、树木、动物都没有了…",sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.45,9.2,0.5,"只剩下钱和金子 💰 — 会发生什么?",sz=22,b=True,c=ALERT,a=PP_ALIGN.CENTER)
tb(s,0.4,2.05,9.2,0.3,"If water, air, trees & animals are all gone — only money & gold left. What happens?",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
for i in range(4):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.45),Inches(2.20),Inches(1.75))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ECO;sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.65,2.1,0.9,"❓",sz=64,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.60,2.1,0.4,f"想法 {i+1}",sz=15,b=True,c=ECO,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.35,"钱能喝吗? 金子能呼吸吗?","Think-Pair-Share: 我觉得会 ___ 。")
n+=1;pn(s,n)
notes(s,"think-pair-share (3 分钟):\n• 不要给答案!\n• 让学生 think → pair → share, 把答案写在白板上。\n• 引导词: 钱不能喝, 金子不能呼吸, 也不能吃…\n• 说完再翻到下一页揭晓。")

# 6. STORY — reveal: what we truly need
s=ns();bg(s,CREAM);hb(s,"💚 我们真正需要的  What We Really Need",ECO)
tb(s,0.4,0.85,9.2,0.4,"钱买不到这些! 它们才是宝贝。",sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,0.4,1.25,9.2,0.3,"Money can't buy these — they are the real treasure.",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
need=[("🌊","水","We drink it",REDUCE),
      ("🌬️","空气","We breathe it",SKY),
      ("🌳","树木","They make air",ECO),
      ("🐾","动物","Our friends",BROWN)]
for i,(em,cn,en,cl) in enumerate(need):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.70),Inches(2.20),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,1.90,2.1,1.0,em,sz=64,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.00,2.1,0.5,cn,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,2.1,0.4,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,"所以我们要保护地球, 还是保护钱?","一起说: 保护地球! 保护水、空气、树、动物!")
n+=1;pn(s,n)
notes(s,"答案揭晓 (2 分钟):\n• 「你们说得对 — 钱不能喝, 不能呼吸!」\n• 强调: 水/空气/树木/动物没有了, 有再多钱也活不下去。\n• 过渡: 「那我们怎么保护它们? 从少制造垃圾开始 — 这就叫 Zero Waste!」")

# 7. WHAT IS ZERO WASTE (8 min · part 2)
s=ns();bg(s,CREAM);hb(s,"🌱 什么是 Zero Waste?  What is Zero Waste?",ECO)
# Contrast: NOT vs IS
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=ALERT;sh.line.width=Pt(2.5)
tb(s,0.55,1.15,4.2,0.5,"❌ 不是…",sz=20,b=True,c=ALERT)
tb(s,0.55,1.75,4.2,1.2,"不是等垃圾多了\n再去处理垃圾。",sz=20,b=True,c=DARK)
tb(s,0.55,3.00,4.2,0.4,"NOT just dealing with trash",sz=12,c=GRAY)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.0),Inches(4.5),Inches(2.5))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=ECO;sh2.line.width=Pt(2.5)
tb(s,5.25,1.15,4.2,0.5,"✅ 而是…",sz=20,b=True,c=ECO)
tb(s,5.25,1.75,4.2,1.2,"努力少制造\n垃圾! 💚",sz=24,b=True,c=ECO)
tb(s,5.25,3.00,4.2,0.4,"Make LESS trash in the first place",sz=12,c=GRAY)
tb(s,0.4,3.70,9.2,0.5,"少一点垃圾 = 帮地球一个大忙!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,"我们怎样才能少制造垃圾?","一起想: 我可以 ___ , 少一点垃圾!")
n+=1;pn(s,n)
notes(s,"讲解 (3 分钟):\n• 关键澄清: Zero Waste 不是「多扔多回收」, 而是「一开始就少产生垃圾」。\n• 打比方: 水龙头一直流水, 光擦地不如先关小水龙头。\n• 过渡: 「有 3 个好办法 — 我们叫它 3R!」")

# 8. INTRO 3R — overview
s=ns();bg(s,CREAM);hb(s,"3️⃣ 认识 3R  Meet the 3R",ECO)
tb(s,0.4,0.85,9.2,0.4,"3 个好办法, 都是 R 开头!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
threeR=[("R","Reduce","减少","用得少一点\nUse less",REDUCE,"⬇️"),
        ("R","Reuse","重复使用","再用一次\nUse again",REUSE,"🔁"),
        ("R","Recycle","回收","变成新东西\nMake new",RECYCLE,"♻️")]
for i,(L,en,cn,d,cl,em) in enumerate(threeR):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.35),Inches(3.0),Inches(3.1))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    pill(s,x+1.15,1.50,0.7,0.5,L,cl,sz=22)
    tb(s,x+0.1,2.10,2.8,0.7,em,sz=54,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.95,2.8,0.4,en,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.40,2.8,0.4,cn,sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tb(s,x+0.1,3.90,2.8,0.4,ls[1],sz=11,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.55,"3R 是哪 3 个? 一起念!","Reduce 减少! Reuse 重复使用! Recycle 回收!")
n+=1;pn(s,n)
notes(s,"引入 3R (2 分钟):\n• 3 个都是 R 开头, 好记!\n• 顺序有意义: 先 Reduce (最好), 再 Reuse, 最后才 Recycle。\n• 接下来 3 页各讲一个 R + 生活例子。")

# 9-11. Each R with examples
s=r_example_slide("R","Reduce","减少","jiǎn shǎo","少用 = 从源头少制造垃圾",REDUCE,
    [("💧","关小水龙头"),("📄","纸用两面"),("🛍️","自带布袋"),("🥤","不用一次性")],
    "💚 最好的办法 — 一开始就用得少!  Best R: use less!");n+=1;pn(s,n)
notes(s,"Reduce (2 分钟):\n• 强调这是「最好的 R」— 没产生的垃圾才是真正的零垃圾。\n• 让学生举例: 家里可以少用什么?")

s=r_example_slide("R","Reuse","重复使用","chóng fù shǐ yòng","同一样东西, 再用一次、很多次",REUSE,
    [("🍶","瓶子当花盆"),("🎒","旧衣服传给弟妹"),("📦","纸箱做玩具"),("🥫","罐子做笔筒")],
    "🔁 别急着扔 — 想想还能做什么!  Don't toss — reuse it!");n+=1;pn(s,n)
notes(s,"Reuse (2 分钟):\n• 「扔之前先想: 它还能做什么?」\n• 这一条正好连到 Session 3 的「变废为宝」项目。")

s=r_example_slide("R","Recycle","回收","huí shōu","扔进回收箱, 工厂把它变成新东西",RECYCLE,
    [("🧴","塑料"),("📰","纸"),("🥫","金属罐"),("🍾","玻璃")],
    "♻️ 分对箱子, 它就有第二次生命!  Sort it right — new life!");n+=1;pn(s,n)
notes(s,"Recycle (2 分钟):\n• 回收要分类: 塑料 / 纸 / 金属 / 玻璃。\n• 提醒: 脏的、混在一起的回收不了 — 所以 Reduce/Reuse 更好。")

# 12. 3R SORTING GAME — how to play (10 min · part 3)
s=ns();bg(s,CREAM);hb(s,"🎮 3R 分类游戏  3R Sorting Game",REUSE)
tb(s,0.4,0.85,9.2,0.4,"上台! 拿一张图片, 放进对的 R 区, 说理由!",sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
zones=[("⬇️ Reduce","减少",REDUCE),("🔁 Reuse","重复使用",REUSE),("♻️ Recycle","回收",RECYCLE)]
for i,(en,cn,cl) in enumerate(zones):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.40),Inches(3.0),Inches(1.5))
    sh.fill.solid();sh.fill.fore_color.rgb=cl;sh.line.fill.background()
    tb(s,x+0.1,1.60,2.8,0.5,en,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.15,2.8,0.5,cn,sz=20,b=True,c=WHITE,a=PP_ALIGN.CENTER)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.10),Inches(9.2),Inches(1.35))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=REUSE;sh3.line.width=Pt(2)
tb(s,0.55,3.20,9.0,0.35,"🎯 怎么玩  How to Play",sz=14,b=True,c=REUSE)
tf=tb(s,0.6,3.60,4.5,0.4,"1️⃣ 老师在墙上贴 3 个 R 区",sz=13,b=True,c=DARK)
ap(tf,"2️⃣ 学生上台抽一张图片",sz=13,b=True,c=DARK)
tf2=tb(s,5.2,3.60,4.4,0.4,"3️⃣ 走到对的 R 区贴上",sz=13,b=True,c=DARK)
ap(tf2,"4️⃣ 说: 我把 ___ 放在 ___ , 因为 ___ 。",sz=12,b=True,c=ECO)
teacher_student_bar(s,4.55,"这张图片应该放哪个 R?","「我把 ___ 放在 ___ , 因为 ___ 。」")
n+=1;pn(s,n)
notes(s,"分类游戏 (10 分钟):\n• 老师准备: 3 张 R 区标牌 (贴墙/地上) + 一叠图片卡。\n• 图片卡例子见下一页。\n• 每个学生上台 1 张, 边贴边说理由。\n• 有的图片可以放多个 R 区 — 允许! 让学生解释就好 (例: 塑料瓶可 Reuse 也可 Recycle)。\n• 低 prep: 图片可从旧杂志剪, 或直接用实物 (瓶子/纸/罐子)。")

# 13. 3R SORTING — sample cards + answers
s=ns();bg(s,CREAM);hb(s,"🃏 分类卡片 + 参考答案  Sorting Cards + Answers",REUSE)
cards=[("💧 一直流的水龙头","Reduce 减少",REDUCE),
       ("📄 只写一面的纸","Reduce 减少",REDUCE),
       ("🍶 空塑料瓶","Reuse / Recycle",REUSE),
       ("👕 小了的衣服","Reuse 重复使用",REUSE),
       ("📦 快递纸箱","Reuse / Recycle",RECYCLE),
       ("🥫 喝完的易拉罐","Recycle 回收",RECYCLE),
       ("📰 旧报纸","Recycle 回收",RECYCLE),
       ("🛍️ 塑料袋","Reduce 减少",REDUCE)]
for i,(item,ans,cl) in enumerate(cards):
    col=i%2;row=i//2
    x=0.4+col*4.75;y=1.05+row*0.98
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(4.55),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,x+0.15,y+0.10,2.7,0.6,item,sz=14,b=True,c=DARK)
    pill(s,x+2.95,y+0.22,1.45,0.42,ans,cl,sz=10)
n+=1;pn(s,n)
notes(s,"参考答案 (老师用):\n• 有些答案不唯一 — 关键是学生能说出理由。\n• 塑料瓶: 洗净可 Reuse (花盆/笔筒), 也可 Recycle。\n• 纸箱: 可 Reuse (做手工) 或 Recycle。\n• 引导「优先级」: 能 Reduce 就 Reduce, 不能才 Reuse, 最后 Recycle。")

# 14. ZERO WASTE DETECTIVE (8 min · part 4)
s=ns();bg(s,CREAM);hb(s,"🕵️ Zero Waste 大侦探  Zero Waste Detective",REDUCE)
tb(s,0.4,0.80,9.2,0.35,"看图找问题 — 说一个更好的做法!",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
cases=[("🚰","水龙头一直开"),("🍽️","一次性餐具"),("💡","没人还开灯"),
       ("📄","纸只用一面"),("🛍️","好多塑料袋"),("🥤","塑料吸管+瓶")]
for i,(em,label) in enumerate(cases):
    col=i%3;row=i//3
    x=0.4+col*3.15;y=1.25+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(1.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=REDUCE;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.10,2.8,0.7,em,sz=40,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,y+0.90,2.8,0.4,f"🔍 {label}",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.55,"侦探! 这里有什么问题? 怎么做更好?","「问题是 ___ 。更好的做法是 ___ 。」")
n+=1;pn(s,n)
notes(s,"大侦探 (8 分钟):\n• 展示 6-10 张「不环保」的图片 (这里 6 个, 可加更多)。\n• 每张: 先找问题, 再说更好的做法。参考答案:\n  - 水龙头一直开 → 用完就关 (Reduce)\n  - 一次性餐具 → 用可洗的餐具 (Reduce/Reuse)\n  - 没人开灯 → 离开关灯 (Reduce)\n  - 纸只用一面 → 两面都用 / 回收 (Reduce/Recycle)\n  - 塑料袋 → 自带布袋 (Reduce/Reuse)\n  - 塑料吸管+瓶 → 自带水杯 (Reduce)\n• 句型: 「问题是___, 更好的做法是___。」")

# 15. SCHOOL ZERO WASTE DISCUSSION (6-8 min · part 5)
s=ns();bg(s,CREAM);hb(s,"🏫 我们学校怎样少制造垃圾?  Our School's Zero Waste",ECO)
tb(s,0.4,0.80,9.2,0.35,"讨论: 学校怎样少扔垃圾 / 少用纸?",sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
ideas=[("♻️","分类回收箱"),("📄","打印用两面"),("🥤","自带水瓶"),
       ("🍱","午餐少用一次性"),("💡","离开教室关灯"),("🖊️","旧纸当草稿")]
for i,(em,label) in enumerate(ideas):
    col=i%3;row=i//3
    x=0.4+col*3.15;y=1.25+row*0.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(3.0),Inches(0.72))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=ECO;sh.line.width=Pt(2)
    tb(s,x+0.1,y+0.15,0.7,0.4,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,x+0.85,y+0.18,2.1,0.4,label,sz=14,b=True,c=DARK)
# Commitment box
cb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.05),Inches(9.2),Inches(1.0))
cb.fill.solid();cb.fill.fore_color.rgb=ECO;cb.line.fill.background()
tb(s,0.55,3.15,9.0,0.4,"✅ 我们选 2–3 个办法, 这两天一起做!",sz=17,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.55,3.60,9.0,0.35,"We choose 2–3 to do together for the next 2 days.  老师写在黑板上 ✍️",sz=12,c=WARM,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.60,"你想让全班一起做哪 2-3 个?","举手: 我选 ___ ! 因为 ___ 。")
n+=1;pn(s,n)
notes(s,"学校讨论 (6-8 分钟):\n• 提问: 「我们学校 / 教室怎样少制造垃圾? 少用纸?」\n• 老师把学生的点子写在黑板上 (上面 6 个是例子, 鼓励更多)。\n• 全班投票选 2-3 个「这两天执行」的行动。\n• 把选中的写下来贴在教室 — Day 5 回顾做到没有。")

# NOTE: 「变废为宝」视频 + Project 过渡 与 3 个项目 已移到 Day 5「环保艺术展」
#       → 见 create_day5_ecoart.py。本课 Session 1 到「学校讨论」结束。

# ---- SESSION 1 → DAY 5 hand-off (light transition, no project detail here) ----
s=ns();bg(s,DEEP)
tb(s,1,0.9,8,0.7,"🎨 下一次…  Coming Up",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.9,8,0.6,"Day 5: 环保艺术展  Eco Art Show",sz=24,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1,2.9,8,0.5,"变废为宝 — 用废物做一件作品! ♻️",sz=18,b=True,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,3.6,8,0.4,"Trash → Treasure. Bring clean bottles / boxes / caps next time!",sz=13,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,4.5,8,0.4,"这两天先当环保小卫士: Reduce · Reuse · Recycle!",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"过渡 (1 分钟):\n• 只做预告 — 不讲项目细节 (项目在 Day 5 环保艺术展)。\n• 请学生开始收集干净的废物 (瓶子/纸箱/瓶盖/旧杂志), Day 5 带来。")

# 17. SESSION 2 DIVIDER
s=div("Session 2  下午 40 min","🔤 复习 + 我会认 5 + 我会写「环保」",REUSE,"📖");n+=1;pn(s,n)

# 18. REVIEW GAME — Baamboozle
s=ns();bg(s,CREAM);hb(s,"🎮 复习游戏  Review · Baamboozle",REUSE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(1.30),Inches(7.0),Inches(3.3))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=REUSE;sh.line.width=Pt(2.5)
tb(s,1.5,1.85,7.0,0.6,"🎲",sz=80,a=PP_ALIGN.CENTER)
tb(s,1.5,2.85,7.0,0.5,"老师在这里放 Baamboozle 链接",sz=20,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,1.5,3.40,7.0,0.4,"Teacher: paste your Baamboozle link here",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1.5,3.90,7.0,0.4,"🔗 https://www.baamboozle.com/...",sz=14,b=True,c=REUSE,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.80,"我们来玩复习游戏!","分组抢答 — 答对得分!")
n+=1;pn(s,n)
notes(s,"复习游戏 (8-10 分钟):\n• 题目涵盖: 什么是 Zero Waste / 3R 各是什么 / 分类 / 大侦探场景 / 词语。\n• 把 Baamboozle 链接粘到这一页, 上课直接打开。\n• 全班分 2-3 组抢答。")

# 19. POINT-TO-WORD GAME
s=ns();bg(s,CREAM);hb(s,"🎯 指对字!  Point to the Right Word!",RECYCLE)
tb(s,0.4,0.85,9.2,0.4,"老师说英文, 你指中文字!",sz=18,b=True,c=DARK,a=PP_ALIGN.CENTER)
words=[("环保","Eco",ECO),("垃圾","Trash",BROWN),("回收","Recycle",RECYCLE),("减少","Reduce",REDUCE),("地球","Earth",SKY)]
for i,(w,en,cl) in enumerate(words):
    x=0.4+i*1.92
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.85),Inches(1.80),Inches(2.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.05,2.30,1.7,1.0,w,sz=32,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.55,1.7,0.5,en,sz=14,b=True,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.55,"「Recycle」在哪? 「Earth」在哪?","用手指 — 大声读出来!")
n+=1;pn(s,n)
notes(s,"游戏 (3 分钟):\n• 老师喊英文, 学生指对应中文字 + 大声读。\n• 反过来: 老师喊中文, 学生说英文。")

# 20-24. 我会认 (5)
read_data=[
    ("环保","huán bǎo","Eco / Protect","🌍","我们都是环保小卫士。"),
    ("垃圾","lā jī","Trash","🗑️","垃圾要少制造, 还要分类。"),
    ("回收","huí shōu","Recycle","♻️","瓶子和纸可以回收。"),
    ("减少","jiǎn shǎo","Reduce","⬇️","减少用纸, 保护大树。"),
    ("地球","dì qiú","Earth","🌏","地球是我们共同的家。"),
]
for w,py,en,em,sent in read_data:
    s=word_card_read(w,py,en,em,sent,color=ECO);n+=1;pn(s,n)

# 25-26. 我会写 (环 · 保)
write_data=[
    ("环","huán","ring / eco",8,"8 笔: 王字旁 + 不  ·  Like a ring 环形 ○"),
    ("保","bǎo","protect",9,"9 笔: 亻+ 呆 (口 + 木)  ·  A person 亻 protects 保护"),
]
for w,py,en,sc,hint in write_data:
    s=word_card_write(w,py,en,sc,hint,color=ECO);n+=1;pn(s,n)

# 27. WORD → PHRASE: 环 + 保 = 环保
s=ns();bg(s,CREAM);hb(s,"🧩 环 + 保 = 环保  Put Them Together",ECO)
tb(s,0.4,1.0,2.8,1.2,"环",sz=110,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,3.1,1.5,1.0,0.8,"+",sz=54,b=True,c=DARK,a=PP_ALIGN.CENTER)
tb(s,4.0,1.0,2.8,1.2,"保",sz=110,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,6.7,1.5,1.0,0.8,"=",sz=54,b=True,c=DARK,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7.6),Inches(1.0),Inches(2.1),Inches(1.4))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=ECO;sh.line.width=Pt(2.5)
tb(s,7.6,1.05,2.1,0.9,"环保",sz=48,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,7.6,2.05,2.1,0.3,"huán bǎo",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(2.85),Inches(9.2),Inches(1.6))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=ECO;sh2.line.width=Pt(2)
tb(s,0.6,2.95,9.0,0.4,"💬 用「环保」说一句话  Make a sentence:",sz=15,b=True,c=ECO)
tf=tb(s,0.6,3.40,9.0,0.4,"·  我是环保小卫士。  I am an eco-guardian.",sz=17,b=True,c=DARK)
ap(tf,"·  我们要环保, 少制造垃圾。  We go eco — make less trash.",sz=17,b=True,c=DARK)
teacher_student_bar(s,4.65,"「环保」是什么意思?","跟读, 再用「环保」说一句话!")
n+=1;pn(s,n)
notes(s,"合词 (3 分钟):\n• 环 (环境) + 保 (保护) = 环保 = 保护环境。\n• 让学生用「环保」造句 (口头), 说得好的写在黑板上。")

# NOTE: Session 3 (变废为宝 + 3 个项目) 已移到 Day 5 → create_day5_ecoart.py

# CLOSING — Zero Waste Hero badge
s=ns();bg(s,DEEP)
tb(s,1,0.6,8,0.7,"🏆 任务完成!",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.35,8,0.5,"Mission Complete!",sz=20,c=WARM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.0),Inches(3.0),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.color.rgb=WHITE;sh.line.width=Pt(4)
tb(s,3.5,2.4,3.0,0.6,"♻️",sz=80,a=PP_ALIGN.CENTER)
tb(s,3.5,3.6,3.0,0.5,"环保小卫士",sz=22,b=True,c=DEEP,a=PP_ALIGN.CENTER)
tb(s,3.5,4.1,3.0,0.4,"Zero Waste Hero",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"这两天记得: Reduce · Reuse · Recycle!  Day 5: 环保艺术展 🎨",sz=14,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# === Save ===
out=os.path.join(BASE,"zerowaste_3R.pptx")
prs.save(out)
print(f"Saved {out}  ({n} slides)")
