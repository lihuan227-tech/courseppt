# Shared helpers for 零 废 弃 与 可 持 续 发 展 (Zero Waste & Sustainability) unit
# Modeled on the 玩转创新科技 helpers, with Zero Waste branding.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ===== Unit palette (Forest & Moss for sustainability theme) =====
INK    = RGBColor(0x0F,0x1A,0x3A)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x88)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
CREAM  = RGBColor(0xF6,0xF8,0xEC)   # Soft sage-tinted cream
WARM   = RGBColor(0xFF,0xF3,0xE0)
IMGBG  = RGBColor(0xE8,0xEE,0xE0)
STAR   = RGBColor(0xF5,0xC2,0x42)

# Forest & Moss sustainability palette
EARTH_GREEN  = RGBColor(0x2C,0x5F,0x2D)   # forest — primary
MOSS         = RGBColor(0x6B,0xA0,0x3D)   # bright leaf
DEEP_TEAL    = RGBColor(0x00,0x69,0x70)   # water
EARTH_BROWN  = RGBColor(0x8B,0x5A,0x2B)   # landfill
FIRE_ORANGE  = RGBColor(0xE6,0x5C,0x00)   # incinerator
RECYCLE_BLUE = RGBColor(0x1B,0x6F,0xBA)   # recycling triangle
LEAF_LIGHT   = RGBColor(0xC8,0xE1,0xA2)   # pale moss
SKY          = RGBColor(0x42,0xA5,0xF5)
GOLD_MEDAL   = RGBColor(0xFF,0xC1,0x07)
SILVER_MEDAL = RGBColor(0xB0,0xB0,0xB0)
BRONZE_MEDAL = RGBColor(0xCD,0x7F,0x32)
OK           = RGBColor(0x2E,0x7D,0x32)   # ✅ green
ALERT        = RGBColor(0xC8,0x25,0x3E)   # ❌ red


def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs


def ns(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def tb(s, l, t, w, h, txt, sz=18, b=False, c=DARK, a=None):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if a:
        p.alignment = a
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = 'KaiTi'
    return tf


def bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def hb(s, txt, c=EARTH_GREEN, t=0.15):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(t), Inches(9.4), Inches(0.55))
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    tb(s, 0.4, t+0.03, 9.2, 0.5, txt, sz=20, b=True, c=WHITE)


def pn(s, n):
    # Top-right corner, inside the header bar — white text so it shows on any header color.
    tb(s, 9.20, 0.21, 0.50, 0.34, str(n), sz=12, b=True, c=WHITE, a=PP_ALIGN.RIGHT)


def notes(s, text):
    nf = s.notes_slide.notes_text_frame
    lines = text.split("\n")
    nf.text = lines[0]
    for line in lines[1:]:
        p = nf.add_paragraph()
        p.text = line


def panel(s, l, t, w, h, color, fill=WHITE, lw=2.5):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.color.rgb = color
    p.line.width = Pt(lw)
    return p


def panel_head(s, l, t, w, color, txt, text_color=WHITE, sz=14):
    h = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.50))
    h.fill.solid()
    h.fill.fore_color.rgb = color
    h.line.fill.background()
    tb(s, l+0.15, t+0.07, w-0.3, 0.40, txt, sz=sz, b=True, c=text_color)


def div(prs, title, sub, color, emoji=""):
    """Section divider slide."""
    s = ns(prs)
    bg(s, color)
    tb(s, 0.3, 1.50, 9.4, 1.0, f"{emoji} {title}", sz=44, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 0.3, 2.70, 9.4, 0.4, sub, sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    for x, y in [(0.8, 4.7), (1.8, 4.5), (7.8, 4.5), (8.6, 4.7)]:
        d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(0.35), Inches(0.35))
        d.fill.solid()
        d.fill.fore_color.rgb = STAR
        d.line.fill.background()
    return s


def tianzi_box(s, x, y, size, char, color, pinyin=None, char_sz=72):
    """Tian-zi-ge (田字格) practice cell."""
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(3)
    mx = x + size/2
    my = y + size/2
    lw = 0.015
    v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx-lw/2), Inches(y), Inches(lw), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = LGRAY; v.line.fill.background()
    h = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(my-lw/2), Inches(size), Inches(lw))
    h.fill.solid(); h.fill.fore_color.rgb = LGRAY; h.line.fill.background()
    tb(s, x, y, size, size, char, sz=char_sz, b=True, c=color, a=PP_ALIGN.CENTER)
    if pinyin:
        tb(s, x, y+size+0.05, size, 0.30, pinyin, sz=12, b=True, c=GRAY, a=PP_ALIGN.CENTER)


def cover(prs, day_num, cn_title, en_title, emoji_row, day_color, inquiry_cn, inquiry_en):
    """Cover slide for Zero Waste unit."""
    s = ns(prs)
    bg(s, INK)
    tb(s, 0.5, 0.40, 9.0, 0.50, "♻️  零 废 弃 与 可 持 续 发 展  Zero Waste & Sustainability",
       sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 1.00, 9.0, 0.50, f"Day {day_num} · {cn_title}",
       sz=30, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 1.55, 9.0, 0.40, en_title,
       sz=16, c=LGRAY, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 2.30, 9.0, 0.70, emoji_row,
       sz=44, c=day_color, a=PP_ALIGN.CENTER)
    qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.50), Inches(7.6), Inches(1.55))
    qbox.fill.solid(); qbox.fill.fore_color.rgb = day_color
    qbox.line.color.rgb = STAR; qbox.line.width = Pt(3)
    tb(s, 1.3, 3.60, 7.4, 0.35, "🤔 今 天 的 探 究 问 题  Today's Inquiry",
       sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 1.3, 4.00, 7.4, 0.45, inquiry_cn, sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, 1.3, 4.55, 7.4, 0.35, inquiry_en, sz=11, c=WARM, a=PP_ALIGN.CENTER)
    return s


def learning_goals(prs, header_color, goals):
    """Learning goals slide."""
    s = ns(prs); bg(s, CREAM); hb(s, "🎯 今 天 的 学 习 目 标  Today's Learning Goals", header_color)
    tb(s, 0.4, 0.85, 9.2, 0.32, "上 完 这 节 课, 你 会……  By the end, you'll be able to…",
       sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    for i, (num, cn, en, cl) in enumerate(goals):
        y = 1.30 + i * 0.95
        panel(s, 0.5, y, 9.0, 0.80, cl, fill=WHITE, lw=2.5)
        tb(s, 0.65, y+0.08, 0.6, 0.40, num, sz=20, b=True, c=cl)
        tb(s, 1.30, y+0.08, 8.0, 0.32, cn, sz=13, b=True, c=DARK)
        tb(s, 1.30, y+0.43, 8.0, 0.30, en, sz=10, c=GRAY)
    return s


def photo_slot(s, l, t, w, h, suggestion_cn, suggestion_en, color=None):
    """Placeholder for a real photo — high contrast so text is readable."""
    if color is None:
        color = EARTH_GREEN
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = IMGBG
    box.line.color.rgb = color; box.line.width = Pt(2)
    tb(s, l, t + 0.10, w, 0.28, "📷 真 实 照 片 · Real Photo",
       sz=10, b=True, c=color, a=PP_ALIGN.CENTER)
    mid_y = t + h / 2 - 0.30
    tb(s, l, mid_y, w, 0.32, "👇 在 这 里 插 入", sz=12, b=True, c=color, a=PP_ALIGN.CENTER)
    tb(s, l + 0.10, mid_y + 0.40, w - 0.20, 0.55, suggestion_cn,
       sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, l + 0.10, mid_y + 0.95, w - 0.20, 0.32, suggestion_en,
       sz=9, c=GRAY, a=PP_ALIGN.CENTER)
    return box


def teacher_box(s, l, t, w, say_cn, do_cn, mins, color=None):
    """Teacher instruction strip: 「老师说 / 学生做 / 时间」.
    Width-aware: w >= 5 → 3-column; otherwise stacked rows.
    Returns total height used.
    """
    if color is None:
        color = EARTH_GREEN
    h = 0.95
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = WARM
    box.line.color.rgb = color; box.line.width = Pt(2)
    # Tags
    tb(s, l + 0.15, t + 0.06, 1.40, 0.22, "👩‍🏫 老师 说",
       sz=9, b=True, c=color)
    tb(s, l + 0.15, t + 0.28, w - 0.30, 0.30, say_cn,
       sz=11, b=True, c=DARK)
    tb(s, l + 0.15, t + 0.58, 1.40, 0.22, "👧 学生 做",
       sz=9, b=True, c=color)
    tb(s, l + 0.15, t + 0.72, w - 1.30, 0.22, do_cn,
       sz=10, c=DARK)
    # Time badge (top-right corner)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l + w - 0.95), Inches(t + 0.10), Inches(0.85), Inches(0.32))
    badge.fill.solid(); badge.fill.fore_color.rgb = STAR; badge.line.fill.background()
    tb(s, l + w - 0.95, t + 0.13, 0.85, 0.28, f"⏱️ {mins}",
       sz=10, b=True, c=INK, a=PP_ALIGN.CENTER)
    return h


def vocab_recognize(prs, header_color, big_emoji, cn, pinyin, en, example_cn, example_en, picture_hint):
    """我会认 slide — single word focus (image LEFT, word card RIGHT, example bottom)."""
    s = ns(prs); bg(s, CREAM)
    hb(s, f"👀 我 会 认 · {cn}  I Can Read", header_color)
    # LEFT: big word card
    panel(s, 0.40, 0.95, 4.55, 3.30, header_color, fill=WARM)
    tb(s, 0.50, 1.10, 4.35, 0.80, big_emoji, sz=58, a=PP_ALIGN.CENTER)
    tb(s, 0.50, 1.95, 4.35, 0.70, cn, sz=44, b=True, c=header_color, a=PP_ALIGN.CENTER)
    tb(s, 0.50, 2.75, 4.35, 0.30, f"{pinyin}  ·  {en}", sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.50, 3.20, 4.35, 0.32, "👉 跟 我 读!  Read after me!",
       sz=12, b=True, c=header_color, a=PP_ALIGN.CENTER)
    # RIGHT: picture placeholder
    photo_slot(s, 5.10, 0.95, 4.50, 3.30, picture_hint,
       "Image placeholder", header_color)
    # Bottom: example sentence
    ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.40), Inches(9.20), Inches(0.95))
    ex.fill.solid(); ex.fill.fore_color.rgb = WHITE
    ex.line.color.rgb = header_color; ex.line.width = Pt(2)
    tb(s, 0.55, 4.48, 9.00, 0.32, "📌 例 句  Example", sz=11, b=True, c=header_color)
    tb(s, 0.55, 4.78, 9.00, 0.35, example_cn, sz=14, b=True, c=DARK)
    tb(s, 0.55, 5.13, 9.00, 0.28, example_en, sz=10, c=GRAY)
    return s


def vocab_write(prs, header_color, cn_phrase, en_word, chars):
    """我会写 slide — day4_emergency format: big-char card + 3 步练习 + 田字格 + teacher/student bar.
    chars is list of (char, pinyin, stroke_count_note, mnemonic)."""
    import re
    s = ns(prs); bg(s, CREAM)
    hb(s, f"✏️ 我 会 写 · {cn_phrase}  I Can Write · {en_word}", header_color)
    py = '  '.join(c[1] for c in chars)
    tot = sum(int(re.search(r'\d+', c[2]).group()) for c in chars if re.search(r'\d+', c[2]))
    # LEFT — big-character card
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0), Inches(4.4), Inches(3.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = WARM; sh.line.color.rgb = header_color; sh.line.width = Pt(2.5)
    csz = 120 if len(cn_phrase) == 1 else 96
    tb(s, 0.5, 1.15, 4.2, 1.95, cn_phrase, sz=csz, b=True, c=header_color, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 3.05, 4.2, 0.45, f"{py}  ·  {en_word}", sz=20, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    if tot:
        tb(s, 0.5, 3.55, 4.2, 0.4, f"{tot} 笔 / {tot} strokes", sz=16, b=True, c=header_color, a=PP_ALIGN.CENTER)
    # RIGHT — 3-step panel
    sh2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.0), Inches(4.6), Inches(1.6))
    sh2.fill.solid(); sh2.fill.fore_color.rgb = WHITE; sh2.line.color.rgb = header_color; sh2.line.width = Pt(2)
    tb(s, 5.15, 1.1, 4.4, 0.4, "✏️ 3 步练习  3 Steps", sz=16, b=True, c=header_color)
    tb(s, 5.15, 1.55, 4.4, 0.35, "1️⃣ 看老师写  Watch teacher", sz=13, c=DARK)
    tb(s, 5.15, 1.90, 4.4, 0.35, "2️⃣ 用手指空中写  Air-write", sz=13, c=DARK)
    tb(s, 5.15, 2.25, 4.4, 0.35, "3️⃣ 在田字格写 3 次", sz=13, c=DARK)
    # RIGHT — 4 田字格 cells
    for i in range(4):
        x = 5.0 + i * 1.15
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.85), Inches(1.05), Inches(1.05))
        sq.fill.solid(); sq.fill.fore_color.rgb = WHITE; sq.line.color.rgb = header_color; sq.line.width = Pt(1.5)
        ln1 = s.shapes.add_connector(1, Inches(x), Inches(3.375), Inches(x + 1.05), Inches(3.375)); ln1.line.color.rgb = LGRAY; ln1.line.width = Pt(0.5); ln1.line.dash_style = 2
        ln2 = s.shapes.add_connector(1, Inches(x + 0.525), Inches(2.85), Inches(x + 0.525), Inches(3.9)); ln2.line.color.rgb = LGRAY; ln2.line.width = Pt(0.5); ln2.line.dash_style = 2
    tb(s, 5.0, 3.95, 4.6, 0.3, "在 田 字 格 里 写 3 次 ↓", sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    # BOTTOM — teacher asks / student does
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.5), Inches(9.2), Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = WARM; bar.line.color.rgb = header_color; bar.line.width = Pt(1.5)
    tb(s, 0.55, 4.55, 4.4, 0.24, "👩‍🏫 老师问 Teacher asks:", sz=10, b=True, c=header_color)
    tb(s, 0.55, 4.80, 4.4, 0.26, f"和 我 一 起 写 「{cn_phrase}」", sz=12, b=True, c=DARK)
    tb(s, 5.05, 4.55, 4.5, 0.24, "🧒 学生 Student does:", sz=10, b=True, c=header_color)
    tb(s, 5.05, 4.80, 4.5, 0.26, "看 → 空中写 → 田字格写 3 次", sz=12, b=True, c=DARK)
    return s


def share_close(prs, day_color, frames_cn, frames_en, next_day_cn, next_day_en, next_emoji="🌱"):
    """Share + close slide."""
    s = ns(prs); bg(s, CREAM); hb(s, "🎤 分 享 + 再 见!  Share + Goodbye", day_color)
    fp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(2.05))
    fp.fill.solid(); fp.fill.fore_color.rgb = INK
    fp.line.color.rgb = STAR; fp.line.width = Pt(3)
    tb(s, 0.55, 1.05, 9.00, 0.35, "💬 今 天 我 学 了 — 句 型  Today's Sentence Frames",
       sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
    for i, frame in enumerate(frames_cn):
        tb(s, 0.55, 1.50 + i*0.50, 9.00, 0.45, frame, sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 2.65, 9.00, 0.28, frames_en, sz=10, c=WARM, a=PP_ALIGN.CENTER)
    np = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.20), Inches(9.20), Inches(2.10))
    np.fill.solid(); np.fill.fore_color.rgb = WHITE
    np.line.color.rgb = day_color; np.line.width = Pt(2.5)
    tb(s, 0.55, 3.35, 9.00, 0.35, "🔮 明 天 见  See You Tomorrow", sz=14, b=True, c=day_color, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 3.80, 9.00, 0.50, f"{next_emoji}  {next_day_cn}", sz=18, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 4.30, 9.00, 0.30, next_day_en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, 0.55, 4.75, 9.00, 0.32, "👋 想 一 想: 你 今 天 最 喜 欢 的 是 什 么?",
       sz=12, b=True, c=day_color, a=PP_ALIGN.CENTER)
    return s


# ===== Camp-style helpers ported from create_day2_camp.py =====
import os as _os


def sentence_frame_bar(s, t, frame_cn, frame_en, color=None):
    """Recurring banner with target sentence frames. Placed at vertical t."""
    if color is None:
        color = EARTH_GREEN
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(t), Inches(9.40), Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb = INK
    sh.line.color.rgb = STAR; sh.line.width = Pt(2)
    tb(s, 0.40, t+0.05, 9.20, 0.28, f"💬 句 型: {frame_cn}",
       sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.40, t+0.33, 9.20, 0.22, frame_en,
       sz=9, c=WARM, a=PP_ALIGN.CENTER)


def answer_panels_slide(prs, header_text, color, panels, img_label="📷 真 实 照 片",
                        img_path=None,
                        subtitle="想 一 想 → 看 答 案 — Now check the answer!",
                        frame_cn=None, frame_en=None):
    """Image-format answer slide: photo LEFT, ✅/❌ panels RIGHT.
    - panels: list of dicts {"q": "...", "mark": "✅"|"❌", "lines": [...]}
    """
    s = ns(prs); bg(s, CREAM); hb(s, header_text, color)
    tb(s, 0.4, 0.74, 9.2, 0.26, subtitle, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
    # LEFT — image area
    img_top = 1.04
    img_h_total = 4.30 if not (frame_cn or frame_en) else 3.55
    img_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.30), Inches(img_top), Inches(4.30), Inches(img_h_total))
    img_box.fill.solid(); img_box.fill.fore_color.rgb = IMGBG
    img_box.line.color.rgb = color; img_box.line.width = Pt(2)
    if img_path and _os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(0.40), Inches(img_top+0.10),
            Inches(4.10), Inches(img_h_total-0.20))
    else:
        tb(s, 0.30, img_top+img_h_total/2-0.20, 4.30, 0.40,
           img_label, sz=13, b=True, c=color, a=PP_ALIGN.CENTER)
    # RIGHT — answer panels (same vertical bounds as image)
    px = 4.80; pw = 4.90
    n_panels = len(panels)
    gap = 0.10
    each_h = (img_h_total - gap*(n_panels-1)) / n_panels
    y = img_top
    for p in panels:
        mark = p["mark"]; lines = p["lines"]; q = p.get("q")
        mark_color = OK if mark == "✅" else ALERT
        pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px), Inches(y), Inches(pw), Inches(each_h))
        pb.fill.solid(); pb.fill.fore_color.rgb = WHITE
        pb.line.color.rgb = mark_color; pb.line.width = Pt(2.5)
        tb(s, px+0.08, y+each_h/2-0.25, 0.55, 0.50,
           mark, sz=22, b=True, c=mark_color, a=PP_ALIGN.CENTER)
        text_x = px + 0.72; text_w = pw - 0.82
        bx = s.shapes.add_textbox(Inches(text_x), Inches(y+0.07),
            Inches(text_w), Inches(each_h-0.14))
        tf = bx.text_frame; tf.word_wrap = True
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        para_idx = 0
        if q:
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            try:
                para.space_before = Pt(0); para.space_after = Pt(1); para.line_spacing = 1.0
            except Exception:
                pass
            r = para.add_run(); r.text = f"❓ {q}"
            r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = mark_color; r.font.name = 'KaiTi'
            para_idx = 1
        for i, ln in enumerate(lines):
            para = tf.paragraphs[0] if (i == 0 and para_idx == 0) else tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            try:
                para.space_before = Pt(1); para.space_after = Pt(1); para.line_spacing = 1.0
            except Exception:
                pass
            r = para.add_run(); r.text = ln
            r.font.size = Pt(10); r.font.bold = True
            r.font.color.rgb = DARK; r.font.name = 'KaiTi'
        y += each_h + gap
    if frame_cn:
        sentence_frame_bar(s, 4.70, frame_cn, frame_en or "", color)
    return s


def zone_slide(prs, emoji, name_cn, name_en, color, rules, frame_cn, frame_en,
               img_label="📷 真 实 照 片"):
    """Category intro slide — emoji + big name on LEFT, rules list RIGHT, sentence frame bar."""
    s = ns(prs); bg(s, CREAM); hb(s, f"{emoji}  {name_cn}  ·  {name_en}", color)
    # LEFT — big category card (shorter so frame bar has room)
    panel(s, 0.40, 0.95, 4.30, 3.35, color, fill=WARM, lw=3)
    tb(s, 0.40, 1.10, 4.30, 1.25, emoji, sz=110, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 2.45, 4.30, 0.55, name_cn, sz=30, b=True, c=color, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 3.05, 4.30, 0.25, name_en, sz=12, c=DARK, a=PP_ALIGN.CENTER)
    # RIGHT — rules / examples list (shorter to match)
    panel(s, 5.00, 0.95, 4.70, 3.35, color, fill=WHITE, lw=2.5)
    panel_head(s, 5.00, 0.95, 4.70, color, "✦ 例 子 + 提 醒  Examples + Tips", sz=12)
    for i, (mark, cn_text, en_text) in enumerate(rules[:5]):
        y = 1.55 + i * 0.54
        tb(s, 5.15, y, 0.45, 0.45,
           mark, sz=18, b=True, c=color, a=PP_ALIGN.LEFT)
        tb(s, 5.65, y, 4.00, 0.28,
           cn_text, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
        tb(s, 5.65, y+0.26, 4.00, 0.22,
           en_text, sz=9, c=GRAY, a=PP_ALIGN.LEFT)
    sentence_frame_bar(s, 4.55, frame_cn, frame_en, color)
    return s


def ab3_slide(prs, title_cn, title_en, question_cn, question_en,
              options, color=None):
    """A/B/C choice slide — 3 options across, big card each.
    options: list of 3 tuples (letter, emoji, label_cn, label_en, color_for_card)
    """
    if color is None:
        color = EARTH_GREEN
    s = ns(prs); bg(s, CREAM); hb(s, f"🎯 {title_cn}  ·  {title_en}", color)
    tb(s, 0.4, 0.85, 9.2, 0.35, question_cn, sz=18, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.4, 1.22, 9.2, 0.28, question_en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    ow = 2.85; ogap = 0.20
    ototal = 3*ow + 2*ogap; ostart = (10 - ototal) / 2
    for i, (letter, em, cn, en, cl) in enumerate(options):
        x = ostart + i * (ow + ogap)
        panel(s, x, 1.65, ow, 3.20, cl, fill=WHITE, lw=3)
        # Letter badge
        bd = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.18), Inches(1.80),
            Inches(0.50), Inches(0.50))
        bd.fill.solid(); bd.fill.fore_color.rgb = cl
        bd.line.color.rgb = STAR; bd.line.width = Pt(2)
        tb(s, x+0.18, 1.83, 0.50, 0.45, letter,
           sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
        # Big emoji
        tb(s, x, 1.85, ow, 1.05, em, sz=66, a=PP_ALIGN.CENTER)
        # Title
        tb(s, x, 3.00, ow, 0.50, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
        tb(s, x+0.10, 3.55, ow-0.20, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    # Vote prompt
    vb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.40), Inches(5.05), Inches(9.20), Inches(0.40))
    vb.fill.solid(); vb.fill.fore_color.rgb = color; vb.line.fill.background()
    tb(s, 0.55, 5.10, 9.0, 0.30,
       "🙋 大 声 喊 出 你 的 选 择: A、B 还 是 C?",
       sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
    return s


def video_slide(prs, title_cn, title_en, before_task, after_action,
                video_url, color=None):
    """Video link slide — book/clip embed with pre-watch task."""
    if color is None:
        color = EARTH_GREEN
    s = ns(prs); bg(s, CREAM); hb(s, f"📺 {title_cn}  ·  {title_en}", color)
    # LEFT — video card with clickable play button
    panel(s, 0.40, 0.95, 4.40, 3.95, color, fill=INK, lw=3)
    tb(s, 0.40, 1.20, 4.40, 1.10, "📺", sz=80, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 2.40, 4.40, 0.45, title_cn, sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 2.88, 4.40, 0.30, title_en, sz=11, c=WARM, a=PP_ALIGN.CENTER)
    # Play button (clickable)
    pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.20), Inches(3.50), Inches(2.80), Inches(0.60))
    pb.fill.solid(); pb.fill.fore_color.rgb = FIRE_ORANGE; pb.line.fill.background()
    tb(s, 1.20, 3.60, 2.80, 0.42, "▶️  点 击 播 放  Click to Play",
       sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    pb.click_action.hyperlink.address = video_url
    tb(s, 0.40, 4.30, 4.40, 0.28, video_url, sz=9, c=LGRAY, a=PP_ALIGN.CENTER)
    # RIGHT — pre-watch task
    panel(s, 5.10, 0.95, 4.50, 3.95, FIRE_ORANGE, fill=WHITE, lw=3)
    panel_head(s, 5.10, 0.95, 4.50, FIRE_ORANGE, "🎬 看 之 前 — 想 一 想", sz=13)
    for i, (em, cn, en) in enumerate(before_task[:3]):
        y = 1.60 + i * 1.00
        tb(s, 5.25, y, 0.55, 0.55, em, sz=28, a=PP_ALIGN.LEFT)
        tb(s, 5.85, y+0.08, 3.60, 0.38, cn, sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
        tb(s, 5.85, y+0.48, 3.60, 0.30, en, sz=9, c=GRAY, a=PP_ALIGN.LEFT)
    # Bottom — after-watch action
    ab = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.45))
    ab.fill.solid(); ab.fill.fore_color.rgb = color; ab.line.fill.background()
    tb(s, 0.55, 5.07, 9.0, 0.32, after_action, sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
    return s


def compare_slide(prs, header_text, left, right, frame_cn=None, frame_en=None):
    """Two-column 用得完 vs 用不完 contrast.
    - left / right: dicts with keys
        "tag_cn", "tag_en"   — column tag (e.g. "用 得 完" / "Runs Out")
        "badge"              — short badge text shown in the corner ("✗ 会 用 完" / "✓ 用 不 完")
        "color"              — column accent color
        "emoji"              — emoji row string (e.g. "🛢️ ⛏️" or "☀️ 💨 💧")
        "title_cn","title_en"— column heading (e.g. "化 石 能 源")
        "bullets"            — list of (mark, cn, en) rows
    LEFT uses a dark (INK) header to read as "the old way"; RIGHT uses its own
    color. Conventions match zone_slide / answer_panels_slide (rounded panels,
    panel_head, STAR accents).
    """
    s = ns(prs); bg(s, CREAM); hb(s, header_text, EARTH_GREEN)
    body_top = 1.00
    body_h = 3.55 if (frame_cn or frame_en) else 4.30
    cols = [(0.30, left, INK), (5.05, right, right.get("color", MOSS))]
    for cx, col, head_c in cols:
        cw = 4.65
        panel(s, cx, body_top, cw, body_h, head_c, fill=WHITE, lw=2.5)
        panel_head(s, cx, body_top, cw, head_c,
                   f"{col['tag_cn']}  {col['tag_en']}", sz=13)
        # badge (corner)
        bd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx + cw - 1.95), Inches(body_top + 0.07),
            Inches(1.80), Inches(0.36))
        bd.fill.solid(); bd.fill.fore_color.rgb = STAR; bd.line.fill.background()
        tb(s, cx + cw - 1.95, body_top + 0.10, 1.80, 0.30,
           col["badge"], sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
        # emoji row
        tb(s, cx, body_top + 0.62, cw, 0.85, col["emoji"], sz=52, a=PP_ALIGN.CENTER)
        # column title
        tb(s, cx, body_top + 1.50, cw, 0.45, col["title_cn"],
           sz=22, b=True, c=head_c, a=PP_ALIGN.CENTER)
        tb(s, cx, body_top + 1.95, cw, 0.28, col["title_en"],
           sz=11, c=GRAY, a=PP_ALIGN.CENTER)
        # bullets
        for i, (mark, cn, en) in enumerate(col["bullets"][:3]):
            y = body_top + 2.32 + i * 0.42
            tb(s, cx + 0.20, y, 0.40, 0.36, mark, sz=15, b=True, c=head_c)
            tb(s, cx + 0.62, y, cw - 0.80, 0.28, cn, sz=11, b=True, c=DARK)
    # center "VS" badge
    vs = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(4.66), Inches(body_top + body_h/2 - 0.32),
        Inches(0.66), Inches(0.66))
    vs.fill.solid(); vs.fill.fore_color.rgb = FIRE_ORANGE
    vs.line.color.rgb = WHITE; vs.line.width = Pt(2.5)
    tb(s, 4.66, body_top + body_h/2 - 0.22, 0.66, 0.45,
       "VS", sz=16, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    if frame_cn:
        sentence_frame_bar(s, 4.70, frame_cn, frame_en or "", EARTH_GREEN)
    return s


# ===== Day 3 (海洋小卫士) reference-structure helpers =====
OCEAN_TINTS = [RGBColor(0x1B,0x6F,0xBA), RGBColor(0x00,0x69,0x70), RGBColor(0x2A,0x9D,0x8F),
               RGBColor(0x3A,0x7C,0xA5), RGBColor(0xE6,0x85,0x3A), RGBColor(0xC8,0x25,0x3E)]


def mission_intro_slide(prs, color, header, cards, frame_cn, frame_en,
                        banner_cn="🌊 今 天 你 是 海 洋 小 卫 士!",
                        banner_en="Today you're a Little Ocean Guardian — save water, cut plastic!"):
    """6 preview cards (emoji + short label) + intro banner + sentence frame.
    cards: list of (emoji, cn, en) — up to 6 (pass SHORT cn labels)."""
    s = ns(prs); bg(s, CREAM); hb(s, header, color)
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.90), Inches(9.20), Inches(0.85))
    band.fill.solid(); band.fill.fore_color.rgb = color; band.line.color.rgb = STAR; band.line.width = Pt(2.5)
    tb(s, 0.5, 0.98, 9.0, 0.40, banner_cn, sz=17, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, 0.5, 1.40, 9.0, 0.28, banner_en, sz=10, c=WHITE, a=PP_ALIGN.CENTER)
    cw, gap = 1.46, 0.10
    total = 6*cw + 5*gap; x0 = (10 - total)/2
    for i, (em, cn, en) in enumerate(cards[:6]):
        x = x0 + i*(cw+gap)
        panel(s, x, 2.05, cw, 1.85, OCEAN_TINTS[i % len(OCEAN_TINTS)], fill=WHITE, lw=2.5)
        tb(s, x, 2.18, cw, 0.72, em, sz=38, a=PP_ALIGN.CENTER)
        tb(s, x, 2.92, cw, 0.48, cn, sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
        tb(s, x, 3.44, cw, 0.26, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    sentence_frame_bar(s, 4.55, frame_cn, frame_en, color)
    return s


def observe_think_slide(prs, color, header, photo_cn, photo_en, senses, dangers,
                        media_cn=None, media_url=None):
    """看+想 station slide: photo LEFT, sensory prompts RIGHT, danger cards bottom.
    senses: list of (emoji, cn) up to 3.  dangers: list of (emoji, cn) up to 4."""
    s = ns(prs); bg(s, CREAM); hb(s, header, color)
    photo_slot(s, 0.30, 0.95, 4.35, 3.30, photo_cn, photo_en, color)
    panel(s, 4.80, 0.95, 4.90, 3.30, color, fill=WHITE, lw=2.5)
    panel_head(s, 4.80, 0.95, 4.90, color, "1️⃣ 👀 看 — 你 看 到 / 想 到 什 么?", sz=12)
    for i, (em, cn) in enumerate(senses[:3]):
        y = 1.55 + i*0.68
        tb(s, 4.95, y, 0.55, 0.55, em, sz=24, a=PP_ALIGN.LEFT)
        tb(s, 5.55, y+0.06, 3.95, 0.55, cn, sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
    if media_cn and media_url:
        link = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.90), Inches(3.55), Inches(4.70), Inches(0.60))
        link.fill.solid(); link.fill.fore_color.rgb = INK; link.line.color.rgb = STAR; link.line.width = Pt(1.5)
        tb(s, 4.95, 3.60, 4.6, 0.26, f"▶️ {media_cn}", sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
        tb(s, 4.95, 3.86, 4.6, 0.24, media_url, sz=7.5, c=LGRAY, a=PP_ALIGN.CENTER)
        link.click_action.hyperlink.address = media_url
    dh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(4.38), Inches(9.40), Inches(0.42))
    dh.fill.solid(); dh.fill.fore_color.rgb = ALERT; dh.line.fill.background()
    tb(s, 0.4, 4.44, 9.2, 0.32, "2️⃣ 🤔 想 — 它 到 了 海 里, 会 伤 害 谁?  我 觉 得……",
       sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
    dw, dgap = 2.28, 0.10; dx0 = (10 - (4*dw+3*dgap))/2
    for i, (em, cn) in enumerate(dangers[:4]):
        x = dx0 + i*(dw+dgap)
        panel(s, x, 4.86, dw, 0.62, color, fill=WHITE, lw=2)
        tb(s, x+0.10, 4.96, 0.5, 0.42, em, sz=20, a=PP_ALIGN.LEFT)
        tb(s, x+0.62, 4.98, dw-0.7, 0.40, cn, sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
    return s


def _ab_rows(s, color, rows, reveal=False):
    """3 scenario rows; each: number badge + scenario + A box (left) + B box (right).
    rows: list of (scenario_cn, a_cn, b_cn, correct)  correct in {'A','B'}."""
    top = 1.72; rh = 0.85
    for i, (sc, a_cn, b_cn, correct) in enumerate(rows[:3]):
        y = top + i*rh
        panel(s, 0.35, y, 9.30, rh-0.13, color, fill=WHITE, lw=2)
        bd = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.46), Inches(y+0.12), Inches(0.40), Inches(0.40))
        bd.fill.solid(); bd.fill.fore_color.rgb = color; bd.line.fill.background()
        tb(s, 0.46, y+0.13, 0.40, 0.36, str(i+1), sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
        tb(s, 1.00, y+0.10, 3.00, 0.55, sc, sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
        for label, txt, bx in (("A", a_cn, 4.10), ("B", b_cn, 6.90)):
            good = reveal and (correct == label)
            oc = OK if good else color
            fillc = RGBColor(0xEF,0xF7,0xEE) if good else WHITE
            panel(s, bx, y+0.13, 2.65, rh-0.39, oc, fill=fillc, lw=(3 if good else 1.5))
            tb(s, bx+0.10, y+0.20, 0.55, 0.35, ("✓"+label if good else label), sz=12, b=True, c=oc, a=PP_ALIGN.LEFT)
            tb(s, bx+0.68, y+0.20, 1.90, 0.35, txt, sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)


def judge_ab_slide(prs, color, header, rows, frame_cn, frame_en):
    """判 A/B — 3 scenario rows, vote bar, sentence frame."""
    s = ns(prs); bg(s, CREAM); hb(s, header, color)
    tb(s, 0.4, 0.78, 9.2, 0.30, "⚖️ 判 — 你 觉 得 呢? A 还 是 B?  (先 选, 再 看!)",
       sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    _ab_rows(s, color, rows, reveal=False)
    vb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.38), Inches(9.20), Inches(0.40))
    vb.fill.solid(); vb.fill.fore_color.rgb = color; vb.line.fill.background()
    tb(s, 0.5, 4.43, 9.0, 0.30, "👉 A → 举 左 手 · B → 举 右 手   🤔 3… 2… 1…",
       sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
    sentence_frame_bar(s, 4.90, frame_cn, frame_en, color)
    return s


def reveal_ab_slide(prs, color, header, rows, do_action, frame_cn, frame_en):
    """答案揭晓 — same 3 rows with green ✓ on safe choice + 做 action + frame."""
    s = ns(prs); bg(s, CREAM); hb(s, header, ALERT)
    tb(s, 0.4, 0.78, 9.2, 0.30, "💡 答 案 揭 晓!  绿 色 ✓ = 保 护 海 洋 的 选 择",
       sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
    _ab_rows(s, color, rows, reveal=True)
    ab = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.38), Inches(9.20), Inches(0.40))
    ab.fill.solid(); ab.fill.fore_color.rgb = OK; ab.line.fill.background()
    tb(s, 0.5, 4.43, 9.0, 0.30, f"🙆 做: {do_action}", sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
    sentence_frame_bar(s, 4.90, frame_cn, frame_en, color)
    return s


def cardgrid_slide(prs, color, header, subtitle, items):
    """6-card overview grid (3 cols × 2 rows). items: list of (emoji, cn, en)."""
    s = ns(prs); bg(s, CREAM); hb(s, header, color)
    tb(s, 0.4, 0.78, 9.2, 0.30, subtitle, sz=13, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    cw, ch, gx, gy = 2.95, 1.75, 0.20, 0.18
    x0 = (10 - (3*cw + 2*gx))/2; y0 = 1.22
    for i, (em, cn, en) in enumerate(items[:6]):
        r, c = divmod(i, 3)
        x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
        panel(s, x, y, cw, ch, OCEAN_TINTS[i % len(OCEAN_TINTS)], fill=WHITE, lw=2.5)
        tb(s, x, y+0.18, cw, 0.70, em, sz=42, a=PP_ALIGN.CENTER)
        tb(s, x, y+0.98, cw, 0.42, cn, sz=19, b=True, c=OCEAN_TINTS[i % len(OCEAN_TINTS)], a=PP_ALIGN.CENTER)
        tb(s, x, y+1.42, cw, 0.26, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    return s


def guess_slide(prs, color, header, subtitle, cards, frame_cn, frame_en):
    """我演你猜 — teacher acts a clue, students guess. cards: (emoji, prompt_cn) up to 6."""
    s = ns(prs); bg(s, CREAM); hb(s, header, color)
    tb(s, 0.4, 0.78, 9.2, 0.28, subtitle, sz=12, b=True, c=GRAY, a=PP_ALIGN.CENTER)
    cw, ch, gx, gy = 2.95, 1.35, 0.20, 0.18
    x0 = (10 - (3*cw + 2*gx))/2; y0 = 1.22
    for i, (em, cn) in enumerate(cards[:6]):
        r, c = divmod(i, 3)
        x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
        panel(s, x, y, cw, ch, OCEAN_TINTS[i % len(OCEAN_TINTS)], fill=WHITE, lw=2)
        tb(s, x, y+0.14, cw, 0.55, em, sz=30, a=PP_ALIGN.CENTER)
        tb(s, x+0.12, y+0.72, cw-0.24, 0.55, cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    sentence_frame_bar(s, 4.62, frame_cn, frame_en, color)
    return s
