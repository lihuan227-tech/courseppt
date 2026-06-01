#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零 废 弃 与 可 持 续 发 展 · Day 1: 垃 圾 去 哪 儿 了?  (Rebuild v2)
Follows day2_camp style — material categories 纸/塑料/食物/金属.
3-session deck for K-5 Chinese immersion summer camp.

Per spec: docs/superpowers/specs/2026-05-31-zero-waste-day1-rebuild-design.md
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = make_presentation()
DAY = EARTH_GREEN
n = 0


def arrow(s, x, y, w=0.20, h=0.30, color=DAY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ============================================================
# DATA — 4 categories + 5 scenarios + vocab + Bamboozle
# ============================================================
CATEGORIES = [
    {
        "emoji": "📜", "name_cn": "纸 类", "name_en": "Paper",
        "color": EARTH_BROWN,
        "examples": [
            ("✅", "报 纸 / 杂 志",  "Newspapers, magazines"),
            ("✅", "纸 箱 / 纸 盒",  "Cardboard boxes"),
            ("✅", "课 本 / 作 业 纸", "Books, worksheets"),
            ("✅", "牛 奶 盒 (洗 干 净)", "Milk cartons (rinse first)"),
            ("⚠️", "脏 餐 巾 不 行!", "Dirty napkins — no!"),
        ],
        "frame_cn": "「这 是 纸 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is paper trash. It can be recycled.",
    },
    {
        "emoji": "🥤", "name_cn": "塑 料", "name_en": "Plastic",
        "color": RECYCLE_BLUE,
        "examples": [
            ("✅", "塑 料 瓶 (水/汽 水)", "Plastic bottles"),
            ("✅", "塑 料 杯 / 盒", "Plastic cups, containers"),
            ("✅", "塑 料 袋 (干 净 的)", "Clean plastic bags"),
            ("✅", "塑 料 玩 具 (坏 了)", "Broken plastic toys"),
            ("⚠️", "塑 料 吸 管 — 看 标 志", "Straws — check the label"),
        ],
        "frame_cn": "「这 是 塑 料 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is plastic trash. It can be recycled.",
    },
    {
        "emoji": "🍌", "name_cn": "食 物 垃 圾", "name_en": "Food Waste",
        "color": MOSS,
        "examples": [
            ("❌", "香 蕉 皮 / 果 核", "Banana peels, fruit cores"),
            ("❌", "剩 饭 / 剩 菜",   "Leftovers"),
            ("❌", "茶 叶 / 咖 啡 渣", "Tea leaves, coffee grounds"),
            ("❌", "鸡 蛋 壳",       "Eggshells"),
            ("💡", "可 以 做 堆 肥!", "Can become compost!"),
        ],
        "frame_cn": "「这 是 食 物 垃 圾. 它 不 可 以 回 收.」",
        "frame_en": "This is food trash. It cannot be recycled.",
    },
    {
        "emoji": "🥫", "name_cn": "金 属", "name_en": "Metal",
        "color": FIRE_ORANGE,
        "examples": [
            ("✅", "易 拉 罐 (汽 水/啤 酒)", "Soda cans"),
            ("✅", "罐 头 盒 (洗 干 净)", "Food cans (rinsed)"),
            ("✅", "金 属 盖 子", "Metal lids"),
            ("✅", "铝 箔 纸 (团 一 团)", "Aluminum foil (balled up)"),
            ("⚠️", "电 池 — 不 能 扔 这 里!", "Batteries — not here!"),
        ],
        "frame_cn": "「这 是 金 属 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is metal trash. It can be recycled.",
    },
]

SCENARIOS = [
    {
        "title": "📰 报 纸 可 以 回 收 吗?",
        "title_en": "Recycle newspaper?",
        "img_label": "📰 一 堆 旧 报 纸",
        "color": EARTH_BROWN,
        "panels": [
            {"mark": "✅", "q": "可 以 回 收 吗?",
             "lines": ["可 以! 它 是 纸 类 垃 圾",
                       "Yes — it's paper",
                       "扔 蓝 色 回 收 桶 · Blue recycle bin"]},
            {"mark": "❌", "q": "扔 厨 房 桶 行 不 行?",
             "lines": ["不 行! 纸 不 是 食 物",
                       "No — paper isn't food",
                       "湿 了 就 不 能 回 收 了 · Wet paper can't be recycled"]},
        ],
        "frame_cn": "「这 是 纸 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is paper trash. It can be recycled.",
    },
    {
        "title": "🥤 塑 料 瓶 可 以 回 收 吗?",
        "title_en": "Recycle plastic bottle?",
        "img_label": "🥤 一 个 透 明 水 瓶",
        "color": RECYCLE_BLUE,
        "panels": [
            {"mark": "✅", "q": "可 以 回 收 吗?",
             "lines": ["可 以! 这 是 塑 料 垃 圾",
                       "Yes — it's plastic",
                       "记 得 倒 空 + 拧 紧 盖 子 · Empty + cap on"]},
            {"mark": "❌", "q": "里 面 有 水 也 可 以 扔 吗?",
             "lines": ["不 行! 要 先 倒 空",
                       "No — pour out water first",
                       "湿 瓶 子 会 污 染 别 的 回 收 物"]},
        ],
        "frame_cn": "「这 是 塑 料 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is plastic trash. It can be recycled.",
    },
    {
        "title": "🍌 香 蕉 皮 可 以 回 收 吗?",
        "title_en": "Recycle banana peel?",
        "img_label": "🍌 一 块 香 蕉 皮",
        "color": MOSS,
        "panels": [
            {"mark": "❌", "q": "可 以 回 收 吗?",
             "lines": ["不 可 以 — 这 是 食 物 垃 圾",
                       "No — it's food waste",
                       "不 能 扔 回 收 桶 · Not in recycling bin"]},
            {"mark": "✅", "q": "那 应 该 去 哪 里?",
             "lines": ["扔 食 物 垃 圾 桶",
                       "Throw in food waste bin",
                       "也 可 以 做 堆 肥! · Or compost!"]},
        ],
        "frame_cn": "「这 是 食 物 垃 圾. 它 不 可 以 回 收.」",
        "frame_en": "This is food trash. It cannot be recycled.",
    },
    {
        "title": "🥛 牛 奶 盒 可 以 回 收 吗?",
        "title_en": "Recycle milk carton?",
        "img_label": "🥛 一 个 空 牛 奶 盒",
        "color": EARTH_BROWN,
        "panels": [
            {"mark": "✅", "q": "可 以 回 收 吗?",
             "lines": ["可 以! 它 是 纸 类",
                       "Yes — it's paper-based",
                       "先 洗 干 净 + 压 扁 · Rinse + flatten"]},
            {"mark": "❌", "q": "里 面 有 牛 奶 也 可 以 吗?",
             "lines": ["不 行! 一 定 要 先 洗",
                       "No — must rinse first",
                       "有 食 物 残 渣 就 不 能 回 收"]},
        ],
        "frame_cn": "「这 是 纸 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is paper trash. It can be recycled.",
    },
    {
        "title": "🥫 易 拉 罐 可 以 回 收 吗?",
        "title_en": "Recycle aluminum can?",
        "img_label": "🥫 一 个 汽 水 易 拉 罐",
        "color": FIRE_ORANGE,
        "panels": [
            {"mark": "✅", "q": "可 以 回 收 吗?",
             "lines": ["可 以! 它 是 金 属 垃 圾",
                       "Yes — it's metal",
                       "金 属 可 以 无 限 次 回 收! · Endless recycling"]},
            {"mark": "💡", "q": "你 知 道 吗?",
             "lines": ["1 个 易 拉 罐 60 天 后",
                       "1 can comes back as a new can in 60 days",
                       "就 能 变 成 新 罐 子!"]},
        ],
        "frame_cn": "「这 是 金 属 垃 圾. 它 可 以 回 收.」",
        "frame_en": "This is metal trash. It can be recycled.",
    },
]

VOCAB_RECOGNIZE = [
    ("📜", "纸",    "zhǐ",      "paper",
     "我 用 纸 写 字.",  "I use paper to write.",
     "📸 白 纸 / 报 纸 / 课 本", EARTH_BROWN),
    ("🥤", "塑 料", "sù liào",  "plastic",
     "这 个 瓶 子 是 塑 料.",  "This bottle is plastic.",
     "📸 塑 料 瓶 + 塑 料 杯", RECYCLE_BLUE),
    ("🍌", "食 物", "shí wù",   "food",
     "食 物 不 能 浪 费.",  "Don't waste food.",
     "📸 各 种 食 物 + 香 蕉 皮", MOSS),
    ("🥫", "金 属", "jīn shǔ",  "metal",
     "易 拉 罐 是 金 属.",  "Cans are metal.",
     "📸 易 拉 罐 + 罐 头 盒", FIRE_ORANGE),
    ("♻️", "回 收", "huí shōu", "recycle",
     "塑 料 瓶 可 以 回 收.",  "Plastic bottles can be recycled.",
     "📸 回 收 标 志 + 蓝 桶", DEEP_TEAL),
]

VOCAB_WRITE = [
    ("纸", "paper", EARTH_BROWN, [
        ("纸", "zhǐ", "7 笔 / 7 strokes", "纟 (绞 丝 旁) + 氏 — 古 时 用 丝 做 纸"),
    ]),
    ("垃 圾", "trash", FIRE_ORANGE, [
        ("垃", "lā", "8 笔 / 8 strokes", "土 字 旁 + 立 — 堆 在 土 上"),
        ("圾", "jī", "6 笔 / 6 strokes", "土 字 旁 + 及"),
    ]),
    ("回 收", "recycle", DEEP_TEAL, [
        ("回", "huí", "6 笔 / 6 strokes", "口 套 一 个 小 口 — 「回 来」"),
        ("收", "shōu", "6 笔 / 6 strokes", "攵 (反 文 旁) — 做 事 + 拿 回"),
    ]),
]


# ============================================================
# 1 · COVER
# ============================================================
cover(prs, 1, "垃 圾 去 哪 儿 了?", "Where Does Our Trash Go?",
      "🗑️  🍌  📰  🥤",
      DAY,
      "你 扔 的 垃 圾 最 后 去 了 哪 里?",
      "Where does the trash you throw away really end up?")
n += 1; pn(prs.slides[-1], n)
notes(prs.slides[-1], "📍 Day 1 · 垃 圾 去 哪 儿 了?\n👩‍🏫 老 师 说: 「今 天 我 们 做 小 侦 探 — 找 出 垃 圾 的 秘 密!」\n⏱️ 1 分 钟")


# ============================================================
# 2 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1",
        "🔍 上 午 10:00–10:45 / 11:00–11:45  ·  垃 圾 去 哪 儿 了?",
        DAY, "🗑️"); n += 1; pn(s, n)


# ============================================================
# 3 · LEARNING GOALS
# ============================================================
s = learning_goals(prs, DAY, [
    ("1️⃣", "认 识 4 种 垃 圾 — 纸 / 塑 料 / 食 物 / 金 属",
     "Recognize 4 trash types", MOSS),
    ("2️⃣", "明 白 垃 圾 不 会 凭 空 消 失",
     "Trash doesn't disappear into thin air", DEEP_TEAL),
    ("3️⃣", "知 道 垃 圾 3 个 去 处 — 回 收 / 填 埋 / 焚 烧",
     "3 places trash goes", RECYCLE_BLUE),
    ("4️⃣", "学 会 句 型: 「这 是 ___ 垃 圾」「它 可 以/不 可 以 回 收」",
     "Use sentence frames", EARTH_GREEN),
])
n += 1; pn(s, n)


# ============================================================
# 4 · WARM UP — 透明垃圾袋
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗑️ 热 身  ·  看 看 老 师 手 里 的 袋 子!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.28,
   "看 看 老 师 手 里 这 个 透 明 的 大 袋 子 — 里 面 装 了 什 么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.15, 9.2, 0.24,
   "Teacher holds up a clear trash bag — what's inside?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)
photo_slot(s, 0.40, 1.50, 4.40, 2.90,
   "📸 透 明 垃 圾 袋 + 4 件 垃 圾",
   "Clear bag with bottle/banana/paper/can", DAY)
warm_items = [
    ("🥤", "塑 料 瓶", "Plastic bottle", RECYCLE_BLUE),
    ("🍌", "香 蕉 皮", "Banana peel", MOSS),
    ("📰", "报 纸",   "Newspaper",     EARTH_BROWN),
    ("🥫", "易 拉 罐", "Aluminum can",   FIRE_ORANGE),
]
for i, (em, cn, en, cl) in enumerate(warm_items):
    row = i // 2; col = i % 2
    x = 5.10 + col * 2.30
    y = 1.50 + row * 1.45
    panel(s, x, y, 2.20, 1.35, cl, fill=WHITE, lw=2.5)
    tb(s, x, y+0.10, 2.20, 0.55, em, sz=32, a=PP_ALIGN.CENTER)
    tb(s, x, y+0.70, 2.20, 0.32, cn, sz=14, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, y+1.05, 2.10, 0.22, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
qbox.fill.solid(); qbox.fill.fore_color.rgb = DAY
qbox.line.color.rgb = STAR; qbox.line.width = Pt(2.5)
tb(s, 0.55, 4.62, 9.0, 0.28, "🎤 老 师 问 3 个 问 题:",
   sz=11, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.92, 9.0, 0.42,
   "1. 你 今 天 扔 过 垃 圾 吗?    2. 垃 圾 去 哪 儿 了?    3. 垃 圾 会 消 失 吗?",
   sz=13, b=True, c=WHITE, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "热 身 4-5 分 钟\n👩‍🏫 老师: 拿 真 实 的 透 明 袋 子 + 4 件 干 净 物 品\n💡 不 给 答 案 — 留 悬 念")


# ============================================================
# 5 · THINK-PAIR-SHARE
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🤝 想 一 想 · 说 一 说  ·  Think-Pair-Share", DAY)
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.05))
qb.fill.solid(); qb.fill.fore_color.rgb = DAY
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 1.08, 9.0, 0.45, "❓ 垃 圾 最 后 去 了 哪 里?",
   sz=24, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.55, 9.0, 0.35, "Where does trash finally end up?",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)
steps = [
    ("🧠", "想 一 想", "Think · 自 己 想 30 秒", DEEP_TEAL),
    ("👥", "两 人 说", "Pair · 跟 同 桌 说 1 分 钟", MOSS),
    ("🎤", "全 班 听", "Share · 几 位 同 学 分 享", FIRE_ORANGE),
]
sw = 2.85; sgap = 0.20
stotal = 3*sw + 2*sgap; sstart = (10 - stotal)/2
for i, (em, cn, en, cl) in enumerate(steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 2.25, sw, 1.85, cl, fill=WHITE, lw=2.5)
    tb(s, x, 2.40, sw, 0.70, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 3.15, sw, 0.38, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.58, sw-0.20, 0.42, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.05))
tbox.fill.solid(); tbox.fill.fore_color.rgb = STAR
tbox.line.color.rgb = DAY; tbox.line.width = Pt(3)
tb(s, 0.55, 4.42, 9.0, 0.40, "⏱️ 计 时: 30 秒 · 60 秒 · 60 秒",
   sz=16, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.85, 9.0, 0.35,
   "Timer: 30s think · 60s pair · 60s share  (~3 min)",
   sz=11, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Think-Pair-Share 3-4 分 钟\n💡 老师 在 黑板 列 出 学生 答案, 不 评 价")


# ============================================================
# 6 · BOOK INTRO — 《垃圾哪里去了》 (video_slide)
# ============================================================
s = video_slide(prs,
    "📖 绘 本 · 《垃 圾 哪 里 去 了?》",
    "Storybook video",
    [("👀", "你 看 到 了 什 么 垃 圾?",  "What trash do you see?"),
     ("🤔", "你 觉 得 垃 圾 会 去 哪?",   "Where do you think trash goes?"),
     ("👂", "认 真 听 + 边 看 边 想!",   "Watch + listen carefully!")],
    "⏱️ 5-7 分 钟 · 老 师 暂 停 4 次 — 后 面 4 页 有 暂 停 点",
    "https://www.youtube.com/watch?v=AKoKVQb9_80",
    color=DAY)
n += 1; pn(s, n)
notes(s, "5-7 分 钟 (含 暂 停 讨 论)\n💡 视 频: youtube.com/watch?v=AKoKVQb9_80")


# ============================================================
# 7-10 · 4 BOOK PAUSE DISCUSSIONS
# ============================================================
pauses = [
    ("Pause 1", "👀", "你 看 到 了 什 么 垃 圾?",
     "What trash did you see?",
     "举 出 你 看 到 的 — 大 家 一 起 数!",
     "Name what you saw — let's count together!",
     MOSS),
    ("Pause 2", "🚛", "垃 圾 车 在 做 什 么?",
     "What is the garbage truck doing?",
     "它 把 垃 圾 运 走 — 运 到 哪 里 呢?",
     "It carries trash away — but to where?",
     DEEP_TEAL),
    ("Pause 3", "❓", "为 什 么 我 们 需 要 垃 圾 车?",
     "Why do we need garbage trucks?",
     "想 一 想: 没 有 垃 圾 车, 会 怎 么 样?",
     "Think: what if there were no trucks?",
     FIRE_ORANGE),
    ("Pause 4", "💭", "垃 圾 会 消 失 吗?",
     "Will the trash disappear?",
     "「消 失」 = 没 有 了. 真 的 没 有 了 吗?",
     "'Disappear' = gone. Is it really gone?",
     RECYCLE_BLUE),
]
for label, em, cn, en, hint_cn, hint_en, cl in pauses:
    s = ns(prs); bg(s, CREAM); hb(s, f"⏸️ {label} · 暂 停 讨 论", cl)
    pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.50))
    pb.fill.solid(); pb.fill.fore_color.rgb = cl; pb.line.fill.background()
    tb(s, 0.55, 1.02, 9.0, 0.38, f"⏸️  {label}  —  老 师 暂 停 视 频",
       sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
    photo_slot(s, 0.40, 1.65, 4.30, 3.20,
       "📸 视 频 暂 停 截 图", "Video pause snapshot", cl)
    panel(s, 5.00, 1.65, 4.60, 3.20, cl, fill=WHITE, lw=3)
    tb(s, 5.10, 1.80, 4.40, 0.80, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, 5.10, 2.70, 4.40, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, 5.10, 3.20, 4.40, 0.32, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    hint_bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.15), Inches(3.65), Inches(4.30), Inches(1.05))
    hint_bx.fill.solid(); hint_bx.fill.fore_color.rgb = WARM
    hint_bx.line.color.rgb = cl; hint_bx.line.width = Pt(1.5)
    tb(s, 5.25, 3.72, 4.10, 0.30, "💡 想 一 想:",
       sz=10, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, 5.25, 4.00, 4.10, 0.35, hint_cn, sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, 5.25, 4.38, 4.10, 0.25, hint_en, sz=9, c=GRAY, a=PP_ALIGN.LEFT)
    tm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.40))
    tm.fill.solid(); tm.fill.fore_color.rgb = STAR; tm.line.fill.background()
    tb(s, 0.55, 5.06, 9.0, 0.30,
       "⏱️ 30-45 秒 讨 论 → 继 续 播 放",
       sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
    n += 1; pn(s, n)
    notes(s, f"{label} 30-45 秒\n💡 不 评 价 — 让 大 家 都 说")


# ============================================================
# 11 · FLOW CHART — 家 → 桶 → 车 → ?
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔄 垃 圾 真 的 不 见 了 吗?", DAY)
tb(s, 0.4, 0.85, 9.2, 0.32,
   "我 们 一 起 跟 着 垃 圾 走 一 走!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
flow = [
    ("🏠", "家",     "Home",        DAY),
    ("🗑️", "垃 圾 桶", "Trash bin",   EARTH_BROWN),
    ("🚛", "垃 圾 车", "Garbage truck", FIRE_ORANGE),
    ("❓", "?",     "Where now?",  RECYCLE_BLUE),
]
fw = 1.95; fgap = 0.30
ftotal = 4*fw + 3*fgap
fstart = (10 - ftotal)/2
for i, (em, cn, en, cl) in enumerate(flow):
    x = fstart + i*(fw + fgap)
    panel(s, x, 1.50, fw, 2.30, cl, fill=WHITE, lw=3)
    tb(s, x, 1.65, fw, 0.85, em, sz=52, a=PP_ALIGN.CENTER)
    tb(s, x, 2.60, fw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.12, fw-0.10, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + fw + 0.02, 2.55, w=0.26, h=0.30, color=DAY)
qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.10), Inches(9.20), Inches(1.30))
qbox.fill.solid(); qbox.fill.fore_color.rgb = INK
qbox.line.color.rgb = STAR; qbox.line.width = Pt(3)
tb(s, 0.55, 4.25, 9.0, 0.50, "🤔 最 后 一 站 — 到 底 在 哪 里?",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.80, 9.0, 0.32,
   "Where is that LAST stop?  (We'll find out next!)",
   sz=12, b=True, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "流 程 图 2-3 分 钟\n💡 不 给 答 案 — 下 一 页 才 猜")


# ============================================================
# 12 · 猜一猜 A/B/C  (ab3_slide)
# ============================================================
s = ab3_slide(prs,
    "猜 一 猜!", "Guess Where!",
    "垃 圾 可 能 去 哪 里? 选 一 个!",
    "Trash could go to ... Pick one!",
    [("A", "♻️", "回 收 中 心", "Recycling Center", RECYCLE_BLUE),
     ("B", "🏔️", "填 埋 场",   "Landfill",         EARTH_BROWN),
     ("C", "🔥", "焚 烧 厂",   "Incinerator",      FIRE_ORANGE)],
    color=DAY)
n += 1; pn(s, n)
notes(s, "猜 一 猜 2 分 钟\n💡 答 案: 三 个 都 对! 真 实 世 界 三 种 都 在 用")


# ============================================================
# 13 · 3 DISPOSAL METHODS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🌍 垃 圾 的 3 种 去 处  ·  3 Places Trash Goes", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "三 个 都 对! 真 实 世 界 里 — 这 3 种 方 法 都 在 用",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
methods = [
    ("♻️", "回 收", "Recycling",
     "瓶 子 → 新 瓶 子", "Bottle → new bottle",
     "📸 工 人 在 分 拣", "Workers sorting",
     "💚 最 好", RECYCLE_BLUE),
    ("🏔️", "填 埋", "Landfill",
     "挖 大 坑 — 埋 起 来", "Big pits — buried",
     "📸 山 一 样 的 垃 圾 堆", "Mountain of trash",
     "⚠️ 占 地 多", EARTH_BROWN),
    ("🔥", "焚 烧", "Incineration",
     "用 火 烧 — 变 成 灰", "Burned → ash",
     "📸 焚 烧 厂 + 烟 囱", "Incinerator chimney",
     "⚡ 有 污 染", FIRE_ORANGE),
]
mw = 2.90; mgap = 0.15
mtotal = 3*mw + 2*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, line_cn, line_en, photo_cn, photo_en, tag, cl) in enumerate(methods):
    x = mstart + i*(mw + mgap)
    panel(s, x, 1.25, mw, 4.05, cl, fill=WHITE, lw=3)
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.25), Inches(mw), Inches(0.85))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.32, mw, 0.40, em + "  " + cn, sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 1.72, mw-0.10, 0.30, en, sz=10, c=WARM, a=PP_ALIGN.CENTER)
    photo_slot(s, x+0.15, 2.20, mw-0.30, 1.75, photo_cn, photo_en, cl)
    tb(s, x+0.10, 4.05, mw-0.20, 0.35, line_cn, sz=13, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 4.40, mw-0.20, 0.28, line_en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 4.78, mw-0.20, 0.30, tag, sz=11, b=True, c=cl, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "3 种 去 处 3 分 钟 · 每 种 一 句 话")


# ============================================================
# 14-17 · 4 CATEGORY INTROS (zone_slide)
# ============================================================
for cat in CATEGORIES:
    s = zone_slide(prs, cat["emoji"], cat["name_cn"], cat["name_en"],
                   cat["color"], cat["examples"],
                   cat["frame_cn"], cat["frame_en"],
                   img_label=f"📸 {cat['name_cn']} — 真 实 照 片")
    n += 1; pn(s, n)
    notes(s, f"{cat['name_cn']} 2-3 分 钟\n💡 用 句 型 跟 读 3 遍")


# ============================================================
# 18-22 · 5 "可以回收吗?" SCENARIOS (answer_panels_slide)
# ============================================================
for sc in SCENARIOS:
    s = answer_panels_slide(prs,
        f"{sc['title']}  ·  {sc['title_en']}",
        sc["color"],
        sc["panels"],
        img_label=sc["img_label"],
        subtitle="想 一 想 → 看 答 案 — Think, then check the answer!",
        frame_cn=sc["frame_cn"],
        frame_en=sc["frame_en"])
    n += 1; pn(s, n)
    notes(s, f"{sc['title']} 1-2 分 钟\n💡 全 班 用 句 型 说 一 遍")


# ============================================================
# 23 · SORT GAME — 4 zones + 8 cards
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🏆 小 组 游 戏  ·  垃 圾 分 类 大 挑 战!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "桌 子 上 有 4 个 区, 一 起 把 8 张 卡 片 放 到 对 的 区 里!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
zones = [
    ("📜", "纸",        "Paper",  EARTH_BROWN),
    ("🥤", "塑 料",     "Plastic", RECYCLE_BLUE),
    ("🍌", "食 物",     "Food",   MOSS),
    ("🥫", "金 属",     "Metal",  FIRE_ORANGE),
]
zw = 2.20; zgap = 0.15
ztotal = 4*zw + 3*zgap; zstart = (10 - ztotal)/2
for i, (em, cn, en, cl) in enumerate(zones):
    x = zstart + i*(zw + zgap)
    panel(s, x, 1.30, zw, 1.55, cl, fill=WHITE, lw=2.5)
    tb(s, x, 1.42, zw, 0.55, em, sz=34, a=PP_ALIGN.CENTER)
    tb(s, x, 2.05, zw, 0.42, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.48, zw-0.10, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.4, 3.00, 9.2, 0.32, "🎴 卡 片 (每 组 一 套 · 8 张):",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
cards = [
    ("📰", "报 纸"), ("📦", "纸 箱"),
    ("🥤", "塑 料 瓶"), ("🍱", "塑 料 盒"),
    ("🍌", "香 蕉 皮"), ("🍎", "苹 果 核"),
    ("🥫", "易 拉 罐"), ("🥫", "罐 头 盒"),
]
cw = 1.05; cgap = 0.08
ctotal = 8*cw + 7*cgap; cstart = (10 - ctotal)/2
for i, (em, cn) in enumerate(cards):
    x = cstart + i*(cw + cgap)
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(3.40), Inches(cw), Inches(1.20))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = DAY; cd.line.width = Pt(1.5)
    tb(s, x, 3.50, cw, 0.55, em, sz=24, a=PP_ALIGN.CENTER)
    tb(s, x+0.02, 4.05, cw-0.04, 0.40, cn, sz=9, b=True, c=DAY, a=PP_ALIGN.CENTER)
inst = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.80), Inches(9.20), Inches(0.65))
inst.fill.solid(); inst.fill.fore_color.rgb = DAY
inst.line.color.rgb = STAR; inst.line.width = Pt(2.5)
tb(s, 0.55, 4.88, 9.0, 0.32,
   "👥 全 班 分 4-5 组 · 一 起 把 卡 片 放 到 对 的 区!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22,
   "Split into 4-5 groups · sort cards into the 4 zones together",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "分 类 游 戏 5-6 分 钟\n💡 不 跑 来 跑 去 — 桌 上 完 成")


# ============================================================
# 24 · REVEAL + MEDALS
# ============================================================
s = ns(prs); bg(s, INK); hb(s, "🏆 公 布 答 案 + 比 一 比!", FIRE_ORANGE)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "老 师 念 答 案 — 每 组 自 己 核 对, 数 一 数 你 们 对 了 几 张!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.15, 9.2, 0.24,
   "Teacher reads answers · groups check + count their correct cards",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
ak = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(1.50), Inches(9.20), Inches(0.85))
ak.fill.solid(); ak.fill.fore_color.rgb = WHITE
ak.line.color.rgb = STAR; ak.line.width = Pt(2.5)
tb(s, 0.55, 1.58, 9.0, 0.28, "✅ 答 案  Answer Key:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
key = ("📰报纸→纸  📦纸箱→纸  🥤塑料瓶→塑料  🍱塑料盒→塑料  "
       "🍌香蕉皮→食物  🍎苹果核→食物  🥫易拉罐→金属  🥫罐头盒→金属")
tb(s, 0.55, 1.88, 9.0, 0.42, key, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
medals = [
    ("🥇", "第 一 名", "1st Place", GOLD_MEDAL),
    ("🥈", "第 二 名", "2nd Place", SILVER_MEDAL),
    ("🥉", "第 三 名", "3rd Place", BRONZE_MEDAL),
]
mw = 2.85; mgap = 0.20
mtotal = 3*mw + 2*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, cl) in enumerate(medals):
    x = mstart + i*(mw + mgap)
    panel(s, x, 2.55, mw, 1.95, cl, fill=WHITE, lw=3)
    tb(s, x, 2.70, mw, 0.85, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 3.55, mw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 4.05, mw-0.20, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.70), Inches(9.20), Inches(0.75))
cb.fill.solid(); cb.fill.fore_color.rgb = STAR; cb.line.fill.background()
tb(s, 0.55, 4.80, 9.0, 0.35,
   "🎉 🎊  全 班 鼓 掌!  Big applause!  🎉 🎊",
   sz=15, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.20, 9.0, 0.22,
   "每 个 组 都 是 「环 保 小 队」!  Every team is an Eco-Squad!",
   sz=10, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "公 布 答 案 3-4 分 钟\n💡 重 点 不 是 输 赢 — 让 大 家 都 觉 得 自 己 是 环 保 小 队")


# ============================================================
# 25 · EXIT TICKET
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎫 出 门 票  ·  用 句 型 说 一 个!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "请 5-8 位 同 学 — 拿 一 件 东 西, 用 句 型 说 一 句!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.30), Inches(8.80), Inches(1.30))
fr1.fill.solid(); fr1.fill.fore_color.rgb = DAY
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.42, 8.50, 0.35, "1️⃣  句 型 一",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.78, 8.50, 0.55, "「这 是 _______ 垃 圾.」",
   sz=26, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.32, 8.50, 0.25, '"This is ___ trash."',
   sz=11, c=WARM, a=PP_ALIGN.CENTER)
fr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(2.80), Inches(8.80), Inches(1.30))
fr2.fill.solid(); fr2.fill.fore_color.rgb = MOSS
fr2.line.color.rgb = STAR; fr2.line.width = Pt(3)
tb(s, 0.75, 2.92, 8.50, 0.35, "2️⃣  句 型 二",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.28, 8.50, 0.55, "「它 可 以 / 不 可 以 回 收.」",
   sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 3.82, 8.50, 0.25, '"It can / cannot be recycled."',
   sz=11, c=WARM, a=PP_ALIGN.CENTER)
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WARM
ex.line.color.rgb = DAY; ex.line.width = Pt(2)
tb(s, 0.55, 4.40, 9.00, 0.30, "💡 举 个 例 子  Example:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.70, 9.00, 0.32,
   "🥤 「这 是 塑 料 垃 圾. 它 可 以 回 收.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.05, 9.00, 0.30,
   "🍌 「这 是 食 物 垃 圾. 它 不 可 以 回 收.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "Exit Ticket 5-6 分 钟 (5-8 位 学生)")


# ============================================================
# 26 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2",
        "📚 下 午 2:00–2:45  ·  词 汇 + 游 戏 复 习",
        DAY, "🔤"); n += 1; pn(s, n)


# ============================================================
# 27-31 · 我会认 × 5
# ============================================================
for em, cn, py, en, ex_cn, ex_en, hint, cl in VOCAB_RECOGNIZE:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)
    notes(s, f"我 会 认 · {cn} 2-3 分 钟\n💡 跟 读 3 遍 + 造 句")


# ============================================================
# 32 · 词汇配对
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔗 配 对 游 戏  ·  Match the Words", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "把 字 跟 图 连 起 来!  Match each word with its picture!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
left_words = [
    ("纸",    EARTH_BROWN),
    ("塑 料", RECYCLE_BLUE),
    ("食 物", MOSS),
    ("金 属", FIRE_ORANGE),
    ("回 收", DEEP_TEAL),
]
right_pics = [
    ("🥫", "C"), ("📜", "E"), ("🍌", "B"), ("♻️", "D"), ("🥤", "A"),
]
for i, (w, cl) in enumerate(left_words):
    y = 1.30 + i * 0.78
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.60), Inches(y), Inches(3.50), Inches(0.65))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = cl; box.line.width = Pt(2.5)
    tb(s, 0.60, y+0.10, 0.50, 0.45, str(i+1),
       sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, 1.15, y+0.10, 2.85, 0.45, w,
       sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
for i, (em, letter) in enumerate(right_pics):
    y = 1.30 + i * 0.78
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.90), Inches(y), Inches(3.50), Inches(0.65))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DAY; box.line.width = Pt(2.5)
    tb(s, 5.90, y+0.10, 0.50, 0.45, letter,
       sz=18, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, 6.45, y+0.05, 2.85, 0.55, em, sz=28, a=PP_ALIGN.LEFT)
tb(s, 4.15, 3.10, 1.65, 0.40, "🤝 配 对",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
arr = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,
    Inches(4.15), Inches(3.50), Inches(1.65), Inches(0.40))
arr.fill.solid(); arr.fill.fore_color.rgb = STAR; arr.line.fill.background()
ak2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(5.20), Inches(9.20), Inches(0.30))
ak2.fill.solid(); ak2.fill.fore_color.rgb = DAY; ak2.line.fill.background()
tb(s, 0.55, 5.22, 9.0, 0.25,
   "🔑 答 案 (老 师 用): 1-E  2-A  3-B  4-C  5-D",
   sz=10, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "配 对 3-4 分 钟\n💡 全 班 喊 答 案 或 找 学 生 上 来 连")


# ============================================================
# 33 · 拍词卡
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "👋 拍 词 卡!  ·  Slap the Word!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "两 人 一 组 — 老 师 念 词, 谁 先 拍 到 谁 赢!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
slap_cards = [
    ("📜", "纸",    EARTH_BROWN),
    ("🥤", "塑 料", RECYCLE_BLUE),
    ("🍌", "食 物", MOSS),
    ("🥫", "金 属", FIRE_ORANGE),
    ("♻️", "回 收", DEEP_TEAL),
    ("👋", "拍!",   STAR),
]
scw = 2.85; scgap = 0.15
scstart = (10 - 3*scw - 2*scgap) / 2
for i, (em, cn, cl) in enumerate(slap_cards):
    row = i // 3; col = i % 3
    x = scstart + col*(scw + scgap)
    y = 1.30 + row*1.40
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(scw), Inches(1.20))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = cl; cd.line.width = Pt(3)
    tb(s, x+0.15, y+0.20, 0.90, 0.85, em, sz=44, a=PP_ALIGN.LEFT)
    tb(s, x+1.10, y+0.30, scw-1.20, 0.55, cn,
       sz=22, b=True, c=cl, a=PP_ALIGN.LEFT)
how_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.40), Inches(9.20), Inches(1.00))
how_box.fill.solid(); how_box.fill.fore_color.rgb = DAY
how_box.line.color.rgb = STAR; how_box.line.width = Pt(2.5)
tb(s, 0.55, 4.50, 9.00, 0.30, "🎯 怎 么 玩:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.80, 9.00, 0.30,
   "• 两 人 一 组, 桌 上 摆 5 张 词 卡   • 老师 念 词   • 谁 先 拍 中 谁 赢!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.13, 9.00, 0.25,
   "Pairs · 5 cards on table · teacher calls · first slap wins",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "拍 词 卡 4-5 分 钟\n💡 轻 拍! 不 是 打")


# ============================================================
# 34 · 句型练习 + Pair Share
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗣️ 句 型 练 习  ·  Sentence Practice", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "和 同 桌 一 起 — 一 人 拿 一 件 东 西, 用 句 型 说 给 对 方 听!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.20), Inches(8.80), Inches(1.05))
fr1.fill.solid(); fr1.fill.fore_color.rgb = DAY
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.30, 8.50, 0.30, "1️⃣  这 是 ___ 垃 圾.",
   sz=20, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.65, 8.50, 0.35,
   "→  这 是 (纸 / 塑 料 / 食 物 / 金 属) 垃 圾.",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.75, 2.02, 8.50, 0.20, "This is ___ trash.",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
fr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(2.45), Inches(8.80), Inches(1.05))
fr2.fill.solid(); fr2.fill.fore_color.rgb = MOSS
fr2.line.color.rgb = STAR; fr2.line.width = Pt(3)
tb(s, 0.75, 2.55, 8.50, 0.30, "2️⃣  它 可 以 / 不 可 以 回 收.",
   sz=20, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 2.90, 8.50, 0.35, "→  它 (可 以 / 不 可 以) 回 收.",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.27, 8.50, 0.20, "It (can / cannot) be recycled.",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
ps = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.75), Inches(9.20), Inches(1.65))
ps.fill.solid(); ps.fill.fore_color.rgb = WARM
ps.line.color.rgb = DAY; ps.line.width = Pt(2.5)
tb(s, 0.55, 3.85, 9.00, 0.32, "👥 Pair Share  ·  和 同 桌 一 起 练!",
   sz=14, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.22, 9.00, 0.30, "🔹 A 同 学: 「这 是 塑 料 垃 圾.」",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.52, 9.00, 0.30, "🔹 B 同 学: 「它 可 以 回 收!」",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.85, 9.00, 0.30,
   "🔄 换 一 件 东 西 — 互 换 角 色, 再 来 一 次!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.18, 9.00, 0.22,
   "Swap the object · switch roles · repeat!",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "句 型 + Pair Share 5-6 分 钟")


# ============================================================
# 35-37 · 我会写 × 3
# ============================================================
for cn_phrase, en, cl, chars in VOCAB_WRITE:
    s = vocab_write(prs, cl, cn_phrase, en, chars)
    n += 1; pn(s, n)
    notes(s, f"我 会 写 · {cn_phrase} 3-4 分 钟\n💡 K-2 描 红, 3-5 自 己 写")


# ============================================================
# 38 · BAMBOOZLE GAME
# ============================================================
s = ns(prs); bg(s, INK)
for x, y in [(0.5, 0.5), (9.0, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
        Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()
tb(s, 0.3, 0.55, 9.4, 0.45, "🎮 GAME TIME!",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(1.20), Inches(8.8), Inches(2.10))
tt.fill.solid(); tt.fill.fore_color.rgb = DAY
tt.line.color.rgb = STAR; tt.line.width = Pt(4)
tb(s, 0.8, 1.40, 8.4, 0.85, "Bamboozle · 垃 圾 分 类 大 复 习!",
   sz=36, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.30, 8.4, 0.45, "纸 / 塑 料 / 食 物 / 金 属 — 10 道 题!",
   sz=18, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.8, 2.80, 8.4, 0.40,
   "Trash Sorting Review — 10 Questions!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)
tb(s, 0.3, 3.65, 9.4, 0.95, "♻️  📰  🥤  🍌  🥫",
   sz=52, a=PP_ALIGN.CENTER)
hf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.55))
hf.fill.solid(); hf.fill.fore_color.rgb = WHITE
hf.line.color.rgb = STAR; hf.line.width = Pt(2)
tb(s, 0.55, 4.92, 9.0, 0.28,
   "💻 老 师: 用 bamboozle_day1_trash.csv 导 入 Bamboozle 游 戏!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22,
   "Teachers: import the CSV to bamboozle.com",
   sz=9, c=GRAY, a=PP_ALIGN.CENTER)
hf.click_action.hyperlink.address = "https://bamboozle.com"
n += 1; pn(s, n)
notes(s, "Bamboozle 10-15 分 钟\n💡 CSV: Chinese/zero_waste零废弃/bamboozle_day1_trash.csv")


# ============================================================
# 39 · CLOSING REFLECTION
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "💭 今 天 我 学 会 了 ...  Today I learned ...", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "请 几 位 同 学 用 句 型 说 一 句 — 今 天 学 到 什 么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(1.30), Inches(9.20), Inches(2.30))
fr_box.fill.solid(); fr_box.fill.fore_color.rgb = DAY
fr_box.line.color.rgb = STAR; fr_box.line.width = Pt(3)
tb(s, 0.55, 1.45, 9.00, 0.40, "🌟 句 型:",
   sz=14, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 1.90, 9.00, 0.55,
   "「今 天 我 学 会 了 ______ 是 ______ 垃 圾.」",
   sz=20, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.50, 9.00, 0.30,
   "Today I learned ___ is ___ trash.",
   sz=12, c=WARM, a=PP_ALIGN.CENTER)
tb(s, 0.55, 2.95, 9.00, 0.50,
   "「它 ______ 回 收.」",
   sz=20, b=True, c=WHITE, a=PP_ALIGN.CENTER)
ex_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.80), Inches(9.20), Inches(1.55))
ex_box.fill.solid(); ex_box.fill.fore_color.rgb = WARM
ex_box.line.color.rgb = DAY; ex_box.line.width = Pt(2)
tb(s, 0.55, 3.90, 9.00, 0.30, "💡 例 子  Examples:",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.25, 9.00, 0.32,
   "🥤 「今 天 我 学 会 了 塑 料 瓶 是 塑 料 垃 圾. 它 可 以 回 收.」",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.65, 9.00, 0.32,
   "🍌 「今 天 我 学 会 了 香 蕉 皮 是 食 物 垃 圾. 它 不 可 以 回 收.」",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.05, 9.00, 0.28,
   "👍 老 师 给 「拍 拍 手」 鼓 励!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "总 结 5 分 钟 · 4-5 位 同 学")


# ============================================================
# 40 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3",
        "🛠️ 下 午 3:00–4:30  ·  我 是 垃 圾 分 类 专 家!",
        DAY, "♻️"); n += 1; pn(s, n)


# ============================================================
# 41 · PROJECT MENU
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🛠️ 今 天 做 什 么?  ·  Today's Projects", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "选 一 个 或 两 个 都 做! 一 起 来 创 造!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
projects = [
    ("Project 1", "🎡", "垃 圾 分 类 转 盘", "Sorting Spinner",
     "🍽️ 纸 盘 · 📌 Brass fastener · 🖍️ 彩 笔",
     "Paper plate · brass fastener · markers", MOSS),
    ("Project 2", "🔑", "环 保 钥 匙 牌", "Shrinky Dink Keychain",
     "✨ Shrinky Dink · 🖍️ 油 性 笔 · 🔑 钥 匙 环 · 🔥 烤 箱",
     "Shrinky Dink · Sharpies · keyring · oven", RECYCLE_BLUE),
]
pw = 4.40; pgap = 0.20
ptotal = 2*pw + pgap; pstart = (10 - ptotal)/2
for i, (label, em, cn, en, mat_cn, mat_en, cl) in enumerate(projects):
    x = pstart + i*(pw + pgap)
    panel(s, x, 1.30, pw, 3.95, cl, fill=WHITE, lw=3)
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.30), Inches(pw), Inches(0.55))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.40, pw, 0.40, label, sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
    tb(s, x, 1.95, pw, 1.25, em, sz=88, a=PP_ALIGN.CENTER)
    tb(s, x, 3.30, pw, 0.50, cn, sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.82, pw-0.10, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x+0.20), Inches(4.20), Inches(pw-0.40), Inches(0.95))
    mb.fill.solid(); mb.fill.fore_color.rgb = WARM
    mb.line.color.rgb = cl; mb.line.width = Pt(1.5)
    tb(s, x+0.30, 4.25, pw-0.60, 0.28, "🎒 材 料:",
       sz=10, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x+0.30, 4.50, pw-0.60, 0.32, mat_cn, sz=10, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, x+0.30, 4.83, pw-0.60, 0.28, mat_en, sz=8, c=GRAY, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "选 项 1-2 分 钟\n💡 K-2 推 荐 Project 1, 3-5 都 可 以")


# ============================================================
# 42 · PROJECT 1 MATERIALS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎡 Project 1 · 材 料  ·  Materials", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "做 一 个 可 以 转 的 垃 圾 分 类 转 盘!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
p1_mats = [
    ("🍽️", "纸 盘",        "Paper plate",        "1 个 / 1 each"),
    ("📌", "Brass fastener", "(老 师 装)",           "1 个 / 1 each"),
    ("🖍️", "彩 笔 / 蜡 笔",  "Markers / crayons",  "几 种 颜 色"),
    ("✂️", "剪 刀",         "Scissors",            "(老 师 帮 忙)"),
]
mw = 2.20; mgap = 0.15
mtotal = 4*mw + 3*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, qty) in enumerate(p1_mats):
    x = mstart + i*(mw + mgap)
    panel(s, x, 1.40, mw, 3.10, MOSS, fill=WHITE, lw=2.5)
    tb(s, x, 1.55, mw, 1.00, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, mw, 0.40, cn, sz=16, b=True, c=MOSS, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.12, mw-0.10, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.50, mw-0.10, 0.50, qty, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)
nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.70))
nb.fill.solid(); nb.fill.fore_color.rgb = STAR; nb.line.fill.background()
tb(s, 0.55, 4.85, 9.0, 0.30,
   "💡 老 师 提 前 准 备 — 每 人 一 套 + 老 师 用 的 brass fastener 工 具",
   sz=12, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.18, 9.0, 0.22,
   "Teachers prep one kit per student + brass-fastener tools",
   sz=9, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "材 料 1-2 分 钟")


# ============================================================
# 43 · PROJECT 1 STEPS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎡 Project 1 · 4 步  ·  4 Steps", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.28,
   "4 步 做 完 — 一 个 可 以 转 的 分 类 转 盘!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
p1_steps = [
    ("1️⃣", "🍽️", "拿 一 个 纸 盘", "Take a paper plate", MOSS),
    ("2️⃣", "✂️", "分 成 4 个 区",  "Divide into 4 sections", DEEP_TEAL),
    ("3️⃣", "🖍️", "每 区 画 一 种 垃 圾", "Draw a trash type in each", FIRE_ORANGE),
    ("4️⃣", "📌", "中 间 装 指 针",   "Add a brass-fastener pointer", RECYCLE_BLUE),
]
sw = 2.20; sgap = 0.15
sstart = (10 - 4*sw - 3*sgap) / 2
for i, (num, em, cn, en, cl) in enumerate(p1_steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 1.30, sw, 3.10, cl, fill=WHITE, lw=2.5)
    nb2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x+sw/2-0.30), Inches(1.40), Inches(0.60), Inches(0.60))
    nb2.fill.solid(); nb2.fill.fore_color.rgb = cl; nb2.line.fill.background()
    tb(s, x+sw/2-0.30, 1.40, 0.60, 0.55, num,
       sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x, 2.15, sw, 0.80, em, sz=48, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.05, sw-0.10, 0.55, cn,
       sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.62, sw-0.10, 0.42, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + sw + 0.01, 2.55, w=0.13, h=0.30, color=MOSS)
zb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
zb.fill.solid(); zb.fill.fore_color.rgb = MOSS
zb.line.color.rgb = STAR; zb.line.width = Pt(2.5)
tb(s, 0.55, 4.62, 9.0, 0.30, "🎯 4 个 区 — 画 这 些:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.92, 9.0, 0.42,
   "📜 纸  ·  🥤 塑 料  ·  🥫 金 属  ·  🍌 食 物",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "步 骤 4-5 分 钟 · 老 师 示 范")


# ============================================================
# 44 · PROJECT 1 TEACHER DEMO
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎡 Project 1 · 老 师 示 范  ·  Teacher Demo", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "看 老 师 一 步 一 步 做!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
panel(s, 0.40, 1.20, 4.40, 3.95, MOSS, fill=WHITE, lw=3)
tb(s, 0.40, 1.30, 4.40, 0.40, "✨ 成 品 样 子",
   sz=14, b=True, c=MOSS, a=PP_ALIGN.CENTER)
photo_slot(s, 0.55, 1.80, 4.10, 3.20,
   "📸 完 成 的 分 类 转 盘",
   "Finished sorting spinner", MOSS)
# Right — labeled zones diagram (textual)
panel(s, 5.10, 1.20, 4.50, 3.95, MOSS, fill=WHITE, lw=3)
tb(s, 5.10, 1.30, 4.50, 0.40, "🎯 转 盘 4 个 区",
   sz=14, b=True, c=MOSS, a=PP_ALIGN.CENTER)
zones4 = [
    ("📜", "纸 类",   EARTH_BROWN),
    ("🥤", "塑 料",   RECYCLE_BLUE),
    ("🥫", "金 属",   FIRE_ORANGE),
    ("🍌", "食 物",   MOSS),
]
for i, (em, lbl, cl) in enumerate(zones4):
    row = i // 2; col = i % 2
    x = 5.25 + col * 2.10
    y = 1.85 + row * 1.50
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(2.00), Inches(1.30))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = cl; cd.line.width = Pt(2.5)
    tb(s, x, y+0.15, 2.00, 0.60, em, sz=36, a=PP_ALIGN.CENTER)
    tb(s, x, y+0.80, 2.00, 0.40, lbl, sz=14, b=True, c=cl, a=PP_ALIGN.CENTER)
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(5.25), Inches(9.20), Inches(0.30))
tip.fill.solid(); tip.fill.fore_color.rgb = STAR; tip.line.fill.background()
tb(s, 0.55, 5.27, 9.0, 0.25,
   "🌟 中 间 装 指 针 — 转 一 转, 玩 一 玩!",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "示 范 3-4 分 钟")


# ============================================================
# 45 · PROJECT 1 WORK TIME + WALK-AROUND TIPS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "⏱️ 动 手 时 间!  ·  Work Time", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "现 在 你 来 做! 老 师 走 动 帮 忙.",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
panel(s, 0.40, 1.20, 4.60, 3.30, MOSS, fill=WHITE, lw=3)
tb(s, 0.40, 1.40, 4.60, 1.50, "⏰", sz=110, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.95, 4.60, 0.55, "30 分 钟",
   sz=42, b=True, c=MOSS, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.55, 4.60, 0.35, "30 minutes · Build time",
   sz=13, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.95, 4.60, 0.32, "🎵 老 师 放 一 首 30 分 钟 的 音 乐",
   sz=10, b=True, c=MOSS, a=PP_ALIGN.CENTER)
panel(s, 5.20, 1.20, 4.40, 3.30, FIRE_ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.20, 1.20, 4.40, FIRE_ORANGE, "🚶 老 师 走 动 提 醒", sz=13)
walk_tips = [
    "✦ 用 中 文 说: 「这 是 ___ 垃 圾」",
    "✦ 鼓 励 K-2 — 帮 他 们 画 + 写",
    "✦ 3-5 — 多 写 字, 多 说 中 文",
    "✦ 帮 忙 装 指 针 (brass fastener)",
    "✦ 提 醒: 4 个 区 都 要 画 满!",
]
for i, t in enumerate(walk_tips):
    tb(s, 5.35, 1.85 + i*0.50, 4.10, 0.40, t,
       sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
rc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.75))
rc.fill.solid(); rc.fill.fore_color.rgb = STAR; rc.line.fill.background()
tb(s, 0.55, 4.75, 9.0, 0.32,
   "✅ 完 成 检 查: 4 个 区 都 有 图 · 指 针 转 得 动 · 用 中 文 说 出 名 字",
   sz=12, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.22,
   "Done check: 4 zones drawn · pointer spins · can name in Chinese",
   sz=9, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "动 手 30 分 钟\n💡 K-2 需 要 老 师 帮 装 brass fastener")


# ============================================================
# 46 · PROJECT 2 MATERIALS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔑 Project 2 · 材 料  ·  Materials", RECYCLE_BLUE)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "用 Shrinky Dink 做 一 个 「环 保 小 达 人」 钥 匙 牌!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
p2_mats = [
    ("✨", "Shrinky Dink",  "Shrink plastic",  "3\"x3\" 一 块"),
    ("🖍️", "油 性 笔",     "Sharpies",        "多 种 颜 色"),
    ("🔑", "钥 匙 环",     "Keyring",         "1 个 / 1 each"),
    ("🔥", "烤 箱",        "Oven (teacher!)", "350°F · 1-3 min"),
]
mw = 2.20; mgap = 0.15
mtotal = 4*mw + 3*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, qty) in enumerate(p2_mats):
    x = mstart + i*(mw + mgap)
    panel(s, x, 1.40, mw, 3.10, RECYCLE_BLUE, fill=WHITE, lw=2.5)
    tb(s, x, 1.55, mw, 1.00, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, mw, 0.40, cn, sz=16, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.12, mw-0.10, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.50, mw-0.10, 0.50, qty, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)
sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.70))
sb.fill.solid(); sb.fill.fore_color.rgb = FIRE_ORANGE; sb.line.fill.background()
tb(s, 0.55, 4.83, 9.0, 0.30,
   "⚠️ 安 全: 烤 箱 只 有 老 师 用!  Safety: Only teachers use the oven!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.15, 9.0, 0.22,
   "学 生 画 好 + 剪 好 → 交 给 老 师 烘 烤",
   sz=10, b=True, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "材 料 1-2 分 钟\n⚠️ 烤 箱 安 全")


# ============================================================
# 47 · PROJECT 2 STEPS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔑 Project 2 · 4 步  ·  4 Steps", RECYCLE_BLUE)
tb(s, 0.4, 0.85, 9.2, 0.28,
   "4 步 做 完 — 一 个 「环 保 小 达 人」 钥 匙 牌!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
p2_steps = [
    ("1️⃣", "🖍️", "在 Shrinky Dink 上 画", "Draw on Shrinky Dink"),
    ("2️⃣", "✂️", "剪 出 形 状",           "Cut out shape"),
    ("3️⃣", "🔥", "老 师 烤 (3-5 分)",     "Teacher bakes (3-5 min)"),
    ("4️⃣", "🔑", "穿 钥 匙 环 — 完 成!",  "Add keyring — done!"),
]
sw = 2.20; sgap = 0.15
sstart = (10 - 4*sw - 3*sgap) / 2
for i, (num, em, cn, en) in enumerate(p2_steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 1.20, sw, 2.20, RECYCLE_BLUE, fill=WHITE, lw=2.5)
    tb(s, x+0.10, 1.30, 0.55, 0.40, num,
       sz=18, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.LEFT)
    tb(s, x, 1.42, sw, 0.85, em, sz=42, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.35, sw-0.10, 0.45, cn,
       sz=11, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.82, sw-0.10, 0.40, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + sw + 0.01, 2.10, w=0.13, h=0.30, color=RECYCLE_BLUE)
di = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.55), Inches(9.20), Inches(1.05))
di.fill.solid(); di.fill.fore_color.rgb = WARM
di.line.color.rgb = RECYCLE_BLUE; di.line.width = Pt(2)
tb(s, 0.55, 3.62, 9.00, 0.32, "🎨 设 计 灵 感  Design Ideas:",
   sz=12, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.95, 9.00, 0.32,
   "♻️ Recycle Hero · 🌱 环 保 小 达 人 · 📜 纸 · 🥤 塑 料 · 🥫 金 属 · 🍌 食 物",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.28, 9.00, 0.25,
   "Eco-Hero badges · 'Recycle Hero' · 4 category icons",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)
sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.70))
sb.fill.solid(); sb.fill.fore_color.rgb = FIRE_ORANGE; sb.line.fill.background()
tb(s, 0.55, 4.83, 9.0, 0.30,
   "⚠️ 烤 箱 只 有 老 师 用 — 学 生 不 碰!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.15, 9.0, 0.22,
   "Only teachers use the oven — students do not touch!",
   sz=10, b=True, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "步 骤 3-4 分 钟")


# ============================================================
# 48 · PROJECT 2 DESIGN IDEAS
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎨 Project 2 · 设 计 灵 感  ·  Design Ideas", RECYCLE_BLUE)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "想 一 想 — 你 的 钥 匙 牌 上 画 什 么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
ideas = [
    ("♻️", "Recycle Hero", "Recycle Hero badge",  RECYCLE_BLUE),
    ("🌱", "环 保 小 达 人", "Eco Hero (CN)",       MOSS),
    ("📜", "纸 类 标 志",   "Paper icon",           EARTH_BROWN),
    ("🥤", "塑 料 标 志",   "Plastic icon",         RECYCLE_BLUE),
    ("🍌", "食 物 标 志",   "Food icon",            MOSS),
    ("🥫", "金 属 标 志",   "Metal icon",           FIRE_ORANGE),
]
iw = 2.85; igap = 0.15
isx = (10 - 3*iw - 2*igap) / 2
for i, (em, cn, en, cl) in enumerate(ideas):
    row = i // 3; col = i % 3
    x = isx + col*(iw + igap)
    y = 1.30 + row*1.55
    panel(s, x, y, iw, 1.40, cl, fill=WHITE, lw=2.5)
    tb(s, x+0.15, y+0.15, 0.95, 0.95, em, sz=46, a=PP_ALIGN.LEFT)
    tb(s, x+1.20, y+0.25, iw-1.30, 0.45, cn,
       sz=15, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x+1.20, y+0.72, iw-1.30, 0.30, en,
       sz=10, c=GRAY, a=PP_ALIGN.LEFT)
tm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
tm.fill.solid(); tm.fill.fore_color.rgb = RECYCLE_BLUE
tm.line.color.rgb = STAR; tm.line.width = Pt(2.5)
tb(s, 0.55, 4.65, 9.0, 0.32, "💡 提 醒:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.95, 9.0, 0.30,
   "Shrinky Dink 烤 完 会 缩 小 一 半 — 画 大 一 点!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.25, 9.0, 0.22,
   "Shrinky Dink shrinks 50% when baked — draw BIG!",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "设 计 2-3 分 钟")


# ============================================================
# 49 · PROJECT 2 SAFETY + WORK TIME
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "⏱️ Project 2 · 动 手 时 间  ·  Work Time", RECYCLE_BLUE)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "现 在 你 来 做! 老 师 烤 — 学 生 不 碰 烤 箱.",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
panel(s, 0.40, 1.20, 4.60, 3.30, RECYCLE_BLUE, fill=WHITE, lw=3)
tb(s, 0.40, 1.40, 4.60, 1.50, "⏰", sz=110, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.95, 4.60, 0.55, "25 + 5 分",
   sz=42, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.55, 4.60, 0.35,
   "25 min draw/cut · 5 min teacher bakes",
   sz=12, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.95, 4.60, 0.32,
   "🎵 老 师 放 一 首 25 分 钟 音 乐",
   sz=10, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.CENTER)
panel(s, 5.20, 1.20, 4.40, 3.30, FIRE_ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.20, 1.20, 4.40, FIRE_ORANGE, "⚠️ 安 全 规 则", sz=13)
safety_tips = [
    "🔥 烤 箱 只 有 老 师 用",
    "👀 学 生 远 离 烤 箱",
    "🧤 老 师 用 手 套 拿",
    "❄️ 烤 完 等 凉 了 再 给 学 生",
    "✋ 不 要 用 嘴 咬 / 玩",
]
for i, t in enumerate(safety_tips):
    tb(s, 5.35, 1.85 + i*0.50, 4.10, 0.40, t,
       sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
rc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.75))
rc.fill.solid(); rc.fill.fore_color.rgb = STAR; rc.line.fill.background()
tb(s, 0.55, 4.75, 9.0, 0.32,
   "✅ 完 成 检 查: 设 计 有 图 + 字 · 剪 干 净 · 老 师 烤 + 装 钥 匙 环",
   sz=12, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.22,
   "Done check: design done · clean cut · teacher bakes + adds keyring",
   sz=9, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "动 手 30 分 钟 · 老 师 看 紧 烤 箱")


# ============================================================
# 50 · GALLERY WALK
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🖼️ 展 示  ·  Gallery Walk", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "把 你 的 作 品 摆 在 桌 上 — 大 家 一 起 看 + 说!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
gw_steps = [
    ("1️⃣", "🖼️", "作 品 摆 桌 上", "Place work on desks"),
    ("2️⃣", "🚶", "全 班 静 静 走 一 圈", "Walk silently"),
    ("3️⃣", "👏", "喜 欢 的 — 拍 拍 桌!", "Tap desk if you love it"),
    ("4️⃣", "🎤", "分 享 4-5 位", "Share 4-5 favorites"),
]
sw = 2.20; sgap = 0.15
sstart = (10 - 4*sw - 3*sgap) / 2
for i, (num, em, cn, en) in enumerate(gw_steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 1.30, sw, 3.10, DAY, fill=WHITE, lw=2.5)
    nb3 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x+sw/2-0.30), Inches(1.40), Inches(0.60), Inches(0.60))
    nb3.fill.solid(); nb3.fill.fore_color.rgb = DAY; nb3.line.fill.background()
    tb(s, x+sw/2-0.30, 1.40, 0.60, 0.55, num,
       sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x, 2.15, sw, 0.80, em, sz=48, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.05, sw-0.10, 0.55, cn,
       sz=12, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.62, sw-0.10, 0.42, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
tb_bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
tb_bx.fill.solid(); tb_bx.fill.fore_color.rgb = STAR; tb_bx.line.fill.background()
tb(s, 0.55, 4.65, 9.0, 0.32,
   "🌟 静 静 走, 不 大 声 — 像 真 的 画 廊!",
   sz=13, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.05, 9.0, 0.25,
   "Walk silently — like a real art gallery!",
   sz=10, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Gallery Walk 5 分 钟")


# ============================================================
# 51 · SHARING (SPINNER FRAME)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎤 分 享 句 型 (转 盘)  ·  Spinner Share", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "用 句 型 介 绍 你 的 转 盘!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.30), Inches(8.80), Inches(1.30))
fr1.fill.solid(); fr1.fill.fore_color.rgb = MOSS
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.42, 8.50, 0.32, "🎡  转 盘 句 型:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.77, 8.50, 0.45,
   "「这 是 我 的 垃 圾 分 类 转 盘. 我 会 回 收 ___.」",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.28, 8.50, 0.25,
   "This is my sorting spinner. I will recycle ___.",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
ex_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(2.95), Inches(9.20), Inches(2.45))
ex_box.fill.solid(); ex_box.fill.fore_color.rgb = WARM
ex_box.line.color.rgb = MOSS; ex_box.line.width = Pt(2)
tb(s, 0.55, 3.05, 9.00, 0.32, "💡 示 例  Examples:",
   sz=12, b=True, c=MOSS, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.45, 9.00, 0.35,
   "🥤 「这 是 我 的 垃 圾 分 类 转 盘. 我 会 回 收 塑 料 瓶.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.90, 9.00, 0.35,
   "📰 「这 是 我 的 垃 圾 分 类 转 盘. 我 会 回 收 报 纸.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.40, 9.00, 0.35,
   "🥫 「这 是 我 的 垃 圾 分 类 转 盘. 我 会 回 收 易 拉 罐.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.90, 9.00, 0.28,
   "👉 用 转 盘 演 示 — 转 一 转, 指 到 哪 一 类 就 说 哪 一 类!",
   sz=11, b=True, c=MOSS, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "分 享 5 分 钟 (转 盘)")


# ============================================================
# 52 · SHARING (KEYCHAIN FRAME)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎤 分 享 句 型 (钥 匙 牌)  ·  Keychain Share", RECYCLE_BLUE)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "用 句 型 介 绍 你 的 钥 匙 牌!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.30), Inches(8.80), Inches(1.30))
fr2.fill.solid(); fr2.fill.fore_color.rgb = RECYCLE_BLUE
fr2.line.color.rgb = STAR; fr2.line.width = Pt(3)
tb(s, 0.75, 1.42, 8.50, 0.32, "🔑  钥 匙 牌 句 型:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.77, 8.50, 0.45,
   "「这 是 我 的 环 保 钥 匙 牌. 我 要 保 护 地 球!」",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.28, 8.50, 0.25,
   "This is my eco-keychain. I will protect the Earth!",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
ex_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(2.95), Inches(9.20), Inches(2.45))
ex_box.fill.solid(); ex_box.fill.fore_color.rgb = WARM
ex_box.line.color.rgb = RECYCLE_BLUE; ex_box.line.width = Pt(2)
tb(s, 0.55, 3.05, 9.00, 0.32, "💡 示 例  Examples:",
   sz=12, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.45, 9.00, 0.35,
   "♻️ 「这 是 我 的 环 保 钥 匙 牌. 我 是 Recycle Hero!」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.90, 9.00, 0.35,
   "🌱 「这 是 我 的 环 保 钥 匙 牌. 我 要 保 护 地 球!」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.40, 9.00, 0.35,
   "📜 「这 是 我 的 环 保 钥 匙 牌. 我 会 回 收 纸!」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.90, 9.00, 0.28,
   "👉 把 钥 匙 牌 挂 在 书 包 上, 每 天 提 醒 自 己!",
   sz=11, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "分 享 5 分 钟 (钥 匙 牌)")


# ============================================================
# 53 · GROUP PHOTO + CELEBRATION
# ============================================================
s = ns(prs); bg(s, INK)
for x, y in [(0.5, 0.5), (9.0, 0.5), (0.5, 4.85), (9.0, 4.85)]:
    d = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
        Inches(x), Inches(y), Inches(0.40), Inches(0.40))
    d.fill.solid(); d.fill.fore_color.rgb = STAR; d.line.fill.background()
tb(s, 0.3, 0.55, 9.4, 0.55, "📸 全 班 大 合 影!",
   sz=32, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.3, 1.20, 9.4, 0.35,
   "Big Group Photo — Hold up your creations!",
   sz=14, c=WARM, a=PP_ALIGN.CENTER)
tb(s, 0.3, 1.85, 9.4, 1.20, "🎉  📸  🎊",
   sz=80, a=PP_ALIGN.CENTER)
clb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(3.20), Inches(8.80), Inches(1.55))
clb.fill.solid(); clb.fill.fore_color.rgb = DAY
clb.line.color.rgb = STAR; clb.line.width = Pt(3)
tb(s, 0.75, 3.35, 8.50, 0.45,
   "🌟 我 们 都 是 「环 保 小 达 人」!",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.75, 3.85, 8.50, 0.32,
   "We are all Eco Heroes!",
   sz=14, c=WARM, a=PP_ALIGN.CENTER)
tb(s, 0.75, 4.25, 8.50, 0.40,
   "🤝 拿 着 作 品 — 一 起 拍 一 张!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "合 影 5 分 钟\n💡 拿 着 作 品 — 难 忘 的 回 忆!")


# ============================================================
# 54 · SHARE + CLOSE
# ============================================================
s = share_close(prs, DAY,
    frames_cn=["「今 天 我 学 会 了 ______ 是 ______ 垃 圾.」",
               "「我 要 保 护 地 球 — 我 会 ______.」"],
    frames_en="I learned ___ is ___ trash · I will protect the Earth by ___",
    next_day_cn="✨ 明 天: 重 复 使 用 + 减 少 浪 费!",
    next_day_en="Tomorrow: Reuse & Reduce!",
    next_emoji="🌍")
n += 1; pn(s, n)
notes(s, "收 尾 5 分 钟\n🌍 完 美 结 束!")


# ============================================================
# Save
# ============================================================
out = os.path.join(os.path.dirname(__file__), "day1_trash.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
