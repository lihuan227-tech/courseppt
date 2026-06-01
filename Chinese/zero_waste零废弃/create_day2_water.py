#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零废弃 Day 2 (afternoon) — 保护水资源 & 减少塑料污染
Ocean Guardians 海洋小卫士 mission. Recolors + reuses the design conventions of
野外生存与探险wilderness_pbl/create_day2_camp.py (header bars, mission cards,
sentence-frame bars, image placeholders) in an ocean palette.

Editable native-pptx; teacher pastes real photos into the 📷 placeholders later.
Run:  python3 create_day2_water.py   ->  day2_water_plastic.pptx (+ uses *_baamboozle.csv)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

HERE=os.path.dirname(os.path.abspath(__file__))
OUT=os.path.join(HERE,"day2_water_plastic.pptx")

prs=Presentation()
prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
W,H=prs.slide_width,prs.slide_height

# --- Ocean palette ---
DEEP  = RGBColor(0x0B,0x55,0x63)   # deep teal — primary / headers / dividers
OCEAN = RGBColor(0x15,0x65,0xA0)   # ocean blue
SKY   = RGBColor(0x4A,0xA3,0xDF)   # sky / light blue
SEAGRN= RGBColor(0x2E,0x8B,0x7A)   # sea green
SAND  = RGBColor(0xE9,0xD8,0xA6)   # sand
CREAM = RGBColor(0xFB,0xF7,0xEC)   # page bg
CORAL = RGBColor(0xE0,0x63,0x3F)   # alert / hurts the ocean
SUNYEL= RGBColor(0xF5,0xC2,0x42)
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x2C,0x2C,0x2C)
GRAY  = RGBColor(0x88,0x88,0x88)
LGRAY = RGBColor(0xBB,0xBB,0xBB)
WARM  = RGBColor(0xFF,0xF3,0xE0)
IMGBG = RGBColor(0xE6,0xEE,0xF2)
OK    = RGBColor(0x38,0x8E,0x3C)
ALERT = CORAL
FONT='KaiTi'

# ---- core helpers (same conventions as the example builder) ----
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT
    return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT
    return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def box(s,l,t,w,h,fill=WHITE,line=None,lw=2.0,rad=True):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    return sh
def ib(s,l,t,w,h,lb="📷 图片",sub="",line=SKY):
    box(s,l,t,w,h,fill=IMGBG,line=line,lw=2)
    tb(s,l+0.1,t+h/2-0.28,w-0.2,0.4,lb,sz=13,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    if sub: tb(s,l+0.1,t+h/2+0.10,w-0.2,0.3,sub,sz=10,c=LGRAY,a=PP_ALIGN.CENTER)
def hb(s,txt,c=DEEP,t=0.18):
    box(s,0.3,t,9.4,0.55,fill=c)
    tb(s,0.45,t+0.07,9.1,0.45,txt,sz=16,b=True,c=WHITE)
def kicker(s,txt,c=OCEAN):
    tb(s,0.45,0.80,9.1,0.3,txt,sz=12,b=True,c=c,a=PP_ALIGN.CENTER)
def pn(s,n): tb(s,9.05,5.28,0.75,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def pill(s,l,t,w,h,txt,c,tc=WHITE,sz=13):
    box(s,l,t,w,h,fill=c); tb(s,l+0.08,t+h/2-0.20,w-0.16,0.4,txt,sz=sz,b=True,c=tc,a=PP_ALIGN.CENTER)
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt

NP=[0]
def page(s):
    NP[0]+=1; pn(s,NP[0]); return s

# ---- sentence-frame bar (💬 我来说) ----
def frame_bar(s,t,cn,en):
    box(s,0.3,t,9.4,0.62,fill=WARM,line=SUNYEL,lw=2)
    tb(s,0.45,t+0.13,1.6,0.4,"💬 我来说",sz=14,b=True,c=CORAL)
    tb(s,1.95,t+0.07,7.6,0.3,cn,sz=15,b=True,c=DARK)
    tb(s,1.95,t+0.34,7.6,0.26,en,sz=10,c=GRAY)

# ---- divider ----
def divider(title,sub,emoji="",c=DEEP,tag=""):
    s=ns(); bg(s,c)
    if tag: tb(s,1,1.05,8,0.5,tag,sz=16,b=True,c=SUNYEL,a=PP_ALIGN.CENTER)
    tb(s,1,1.85,8,1.2,f"{emoji} {title}",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,3.15,8,0.8,sub,sz=20,c=RGBColor(0xCF,0xE6,0xF0),a=PP_ALIGN.CENTER)
    NP[0]+=1
    return s

# ---- mission card (numbered task) ----
def mission_card(s,l,t,w,h,num,cn,en,emoji,c):
    box(s,l,t,w,h,fill=WHITE,line=c,lw=2.5)
    bd=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l+0.12),Inches(t+0.10),Inches(0.5),Inches(0.5))
    bd.fill.solid(); bd.fill.fore_color.rgb=c; bd.line.fill.background(); bd.shadow.inherit=False
    tb(s,l+0.12,t+0.17,0.5,0.4,str(num),sz=17,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+0.78,w-0.1,0.7,emoji,sz=40,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+1.55,w-0.1,0.4,cn,sz=16,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+1.95,w-0.1,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)

# ---- simple icon card (emoji + cn + en) ----
def icon_card(s,l,t,w,h,emoji,cn,en,c=OCEAN,em_sz=40):
    box(s,l,t,w,h,fill=WHITE,line=c,lw=2)
    tb(s,l,t+0.18,w,0.7,emoji,sz=em_sz,a=PP_ALIGN.CENTER)
    tb(s,l,t+h-1.05,w,0.4,cn,sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
    if en: tb(s,l,t+h-0.62,w,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)

# ---- vocab recognition card (我会认) ----
def vocab_slide(i,total,char,py,en,word,emoji,c):
    s=ns(); bg(s,CREAM); hb(s,f"📖 我会认 {i}/{total} · I can read")
    kicker(s,"认一认这个字词 — 看图、读一读 Recognize the word")
    # left big char card
    box(s,0.4,1.25,4.3,3.55,fill=WHITE,line=c,lw=3)
    tb(s,0.4,1.55,4.3,1.0,emoji,sz=58,a=PP_ALIGN.CENTER)
    tb(s,0.4,2.75,4.3,1.0,char,sz=62,b=True,c=c,a=PP_ALIGN.CENTER)
    tb(s,0.4,3.95,4.3,0.4,py,sz=22,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.4,4.35,4.3,0.4,en,sz=16,b=True,c=DARK,a=PP_ALIGN.CENTER)
    # right: photo placeholder + usage word
    ib(s,5.0,1.25,4.6,2.05,f"📷 {char} 实物照片",f"insert a photo of «{char}»",line=c)
    box(s,5.0,3.45,4.6,1.35,fill=WHITE,line=SKY,lw=2)
    tb(s,5.2,3.58,4.2,0.35,"🗣️ 组词 Make a word",sz=13,b=True,c=OCEAN)
    tb(s,5.0,3.95,4.6,0.7,word,sz=30,b=True,c=c,a=PP_ALIGN.CENTER)
    return page(s)

# ---- 田字格 writing (我会写) ----
def tian(s,l,t,side,ch=""):
    box(s,l,t,side,side,fill=WHITE,line=CORAL,lw=2.5)
    cx=l+side/2; cy=t+side/2
    g1=s.shapes.add_connector(2,Inches(cx),Inches(t),Inches(cx),Inches(t+side))
    g1.line.color.rgb=RGBColor(0xE7,0xB6,0xB6); g1.line.width=Pt(0.75)
    g2=s.shapes.add_connector(2,Inches(l),Inches(cy),Inches(l+side),Inches(cy))
    g2.line.color.rgb=RGBColor(0xE7,0xB6,0xB6); g2.line.width=Pt(0.75)
    if ch: tb(s,l,t+side/2-side*0.42,side,side*0.85,ch,sz=int(side*54),c=RGBColor(0xCF,0xDD,0xD9),a=PP_ALIGN.CENTER)

def write_slide(i,total,word,en,strokes):
    s=ns(); bg(s,CREAM); hb(s,f"✍️ 我会写 {i}/{total} · I can write")
    kicker(s,f"写一写:{word}  ·  先描红,再自己写一个 — trace, then write")
    chars=list(word); cells=chars+['']*len(chars)
    n=len(cells); side=min(2.1,(9.0-0.35*(n-1))/n); gap=0.35
    tot=side*n+gap*(n-1); x0=(10-tot)/2; y0=1.55
    for j,ch in enumerate(cells):
        tian(s,x0+j*(side+gap),y0,side,ch)
    pill(s,(10-6.2)/2,y0+side+0.35,6.2,0.55,f"✏️ 笔画 {strokes} · {en}",SEAGRN,sz=14)
    return page(s)

print("helpers ready")

# ================================================================ content slides
def s_cover():
    s=ns(); bg(s,DEEP)
    # wave band
    box(s,0,4.55,10,1.07,fill=OCEAN,rad=False)
    box(s,0,4.35,10,0.25,fill=SKY,rad=False)
    tb(s,0.5,0.55,9,0.5,"DAY 2 · 下午 AFTERNOON",sz=15,b=True,c=SUNYEL,a=PP_ALIGN.CENTER)
    tb(s,0.5,1.35,9,1.1,"🌊 保护水资源 & 减少塑料污染",sz=38,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.55,9,0.5,"Protect Water  &  Reduce Plastic Pollution",sz=18,c=RGBColor(0xCF,0xE6,0xF0),a=PP_ALIGN.CENTER)
    box(s,2.3,3.25,5.4,0.7,fill=SEAGRN)
    tb(s,2.3,3.40,5.4,0.45,"🛡️ 海洋小卫士大行动 · Ocean Guardians",sz=17,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,4.75,9,0.5,"谷雨中文 GR EDU · 零废弃 Zero Waste · Day 2",sz=12,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    NP[0]+=1

def s_recap():
    s=ns(); bg(s,CREAM); hb(s,"🔁 复习上午 · 我们学了什么? Morning recap")
    kicker(s,"先想一想,再说一说 — Think, then share")
    ib(s,0.4,1.25,4.1,3.35,"📷 上午活动照片 / 板书","morning recap photo")
    box(s,4.75,1.25,4.85,3.35,fill=WHITE,line=OCEAN,lw=2.5)
    tf=tb(s,5.0,1.45,4.4,0.5,"🤔 想一想 Think back:",sz=15,b=True,c=OCEAN)
    for q in ["• 早上我们学了什么?","• 我们做了什么活动?","• 你记得哪个新词?"]:
        ap(tf,q,sz=14,c=DARK)
    ap(tf,"",sz=6)
    ap(tf,"👉 下午:我们当「海洋小卫士」!",sz=15,b=True,c=SEAGRN)
    ap(tf,"This afternoon — we become Ocean Guardians!",sz=11,c=GRAY)
    frame_bar(s,4.78,"早上我学了 ____ 。","This morning I learned ____ .")
    page(s)

def s_mission():
    s=ns(); bg(s,CREAM); hb(s,"🛡️ 今天的任务 · 海洋小卫士 4 个任务 · 4 Missions")
    kicker(s,"完成 4 个任务,成为海洋小卫士! Finish 4 missions to become a Guardian")
    tasks=[(1,"水很重要","Water matters","💧",OCEAN),
           (2,"海洋告急","Ocean in trouble","🌊",CORAL),
           (3,"一次性的麻烦","Single-use trouble","🥤",SUNYEL if False else SEAGRN),
           (4,"我能做到","I can help","💪",DEEP)]
    w=2.18; gap=0.18; x=(10-(w*4+gap*3))/2
    for k,(n,cn,en,em,c) in enumerate(tasks):
        mission_card(s,x+k*(w+gap),1.45,w,2.6,n,cn,en,em,c)
    page(s)

def s_goals():
    s=ns(); bg(s,CREAM); hb(s,"🎯 学习目标 · 今天要学会什么 Today's Goals")
    goals=[("知道水很重要 — 我们都需要水","Water is important — we all need it","💧",OCEAN),
           ("知道塑料会伤害海洋动物","Plastic hurts ocean animals","🐢",CORAL),
           ("明白「一次性塑料」的问题","Single-use plastic is a problem","🥤",SEAGRN),
           ("会说减少塑料的方法","I can name ways to use less plastic","💪",DEEP)]
    y=1.25; gh=0.95
    for i,(cn,en,em,c) in enumerate(goals):
        box(s,0.4,y,9.2,gh-0.12,fill=WHITE,line=c,lw=2)
        bd=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.6),Inches(y+0.16),Inches(0.5),Inches(0.5))
        bd.fill.solid(); bd.fill.fore_color.rgb=c; bd.line.fill.background(); bd.shadow.inherit=False
        tb(s,0.6,y+0.23,0.5,0.4,str(i+1),sz=17,b=True,c=WHITE,a=PP_ALIGN.CENTER)
        tb(s,1.3,y+0.07,7.0,0.45,cn,sz=17,b=True,c=DARK)
        tb(s,1.3,y+0.48,7.0,0.3,en,sz=11,c=GRAY)
        tb(s,8.7,y+0.05,0.8,0.7,em,sz=30,a=PP_ALIGN.CENTER)
        y+=gh
    page(s)

# ---- generic: header + 4 icon cards + frame ----
def cards4_slide(htxt,kick,items,frame=None,head_c=DEEP,card_c=OCEAN):
    s=ns(); bg(s,CREAM); hb(s,htxt,c=head_c)
    if kick: kicker(s,kick)
    w=2.18; gap=0.18; x=(10-(w*4+gap*3))/2
    ch=2.45 if frame else 3.1; y=1.3
    for k,(em,cn,en) in enumerate(items):
        icon_card(s,x+k*(w+gap),y,w,ch,em,cn,en,c=card_c)
    if frame: frame_bar(s,4.75,frame[0],frame[1])
    return page(s)

def s_water_need():
    cards4_slide("💧 任务 1 · 为什么需要水? Why we need water",
        "人、动物、植物 — 都离不开水 Everyone needs water",
        [("🥤","喝水","drink"),("🛁","洗澡洗手","wash"),("🌱","浇植物","plants"),("🐟","动物","animals")],
        head_c=OCEAN,card_c=OCEAN)

def s_water_precious():
    s=ns(); bg(s,CREAM); hb(s,"💧 任务 1 · 地球的水很珍贵 Earth's water is precious",c=OCEAN)
    kicker(s,"大部分是海水(咸),能喝的淡水很少 — most is salty; little is drinkable")
    ib(s,0.4,1.25,4.1,3.0,"📷 地球 / 一滴水 图片","Earth or a water drop")
    box(s,4.75,1.25,4.85,3.0,fill=WHITE,line=OCEAN,lw=2.5)
    tf=tb(s,5.0,1.45,4.4,0.5,"🌍 你知道吗 Did you know?",sz=15,b=True,c=OCEAN)
    for q in ["🌊 地球上大部分水是海水 — 又咸,不能直接喝","💧 能喝的淡水很少很少","🚰 所以,我们要节约用水、保护水!"]:
        ap(tf,q,sz=14,c=DARK); ap(tf,"",sz=4)
    frame_bar(s,4.45,"水很重要,因为 ____ 。","Water is important because ____ .")
    page(s)

def s_plastic_to_sea():
    s=ns(); bg(s,CREAM); hb(s,"🌊 任务 2 · 塑料去了海里 Plastic reaches the sea",c=CORAL)
    kicker(s,"乱丢的塑料 → 风和雨水 → 河 → 大海 The path of littered plastic")
    steps=[("🗑️","乱丢塑料","litter"),("🌧️","雨水冲走","rain"),("🏞️","流进河里","river"),("🌊","到了大海","ocean")]
    w=1.9; gap=0.45; x=(10-(w*4+gap*3))/2; y=1.7
    for k,(em,cn,en) in enumerate(steps):
        icon_card(s,x+k*(w+gap),y,w,2.2,em,cn,en,c=CORAL)
        if k<3: tb(s,x+w+(w+gap)*k-0.05,y+0.75,0.5,0.5,"→",sz=24,b=True,c=CORAL,a=PP_ALIGN.CENTER)
    frame_bar(s,4.45,"塑料不见了吗? 不 — 它去了大海。","Plastic doesn't disappear — it goes to the ocean.")
    page(s)

def s_animals_hurt():
    s=ns(); bg(s,CREAM); hb(s,"🐢 任务 2 · 塑料伤害海洋动物 Plastic hurts animals",c=CORAL)
    kicker(s,"看一看 — 海洋动物怎么了? What happens to the animals?")
    items=[("🐢","海龟把塑料袋当水母吃了","turtle eats a bag"),
           ("🐟","小鱼被塑料圈困住","fish trapped in a ring"),
           ("🐦","海鸟吃了塑料碎片","seabird eats plastic"),
           ("🦭","海豹被渔网缠住","seal tangled in a net")]
    w=2.18; gap=0.18; x=(10-(w*4+gap*3))/2; y=1.3
    for k,(em,cn,en) in enumerate(items):
        box(s,x+k*(w+gap),y,w,2.45,fill=WHITE,line=CORAL,lw=2)
        tb(s,x+k*(w+gap),y+0.16,w,0.7,em,sz=38,a=PP_ALIGN.CENTER)
        tb(s,x+k*(w+gap)+0.1,y+1.05,w-0.2,0.95,cn,sz=12.5,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+k*(w+gap)+0.1,y+2.05,w-0.2,0.32,en,sz=9.5,c=GRAY,a=PP_ALIGN.CENTER)
    frame_bar(s,4.75,"塑料让 ____ 受伤。","Plastic hurts ____ .")
    page(s)

def s_turtle_story():
    s=ns(); bg(s,CREAM); hb(s,"🐢 任务 2 · 一只海龟的故事 A sea turtle's story",c=CORAL)
    kicker(s,"想一想:我们能帮海龟吗? Can we help the turtle?")
    ib(s,0.4,1.25,4.1,3.35,"📷 海龟 / 塑料袋 图片","turtle & plastic bag")
    box(s,4.75,1.25,4.85,3.35,fill=WHITE,line=SEAGRN,lw=2.5)
    tf=tb(s,5.0,1.45,4.4,0.5,"📖 故事 Story",sz=15,b=True,c=SEAGRN)
    for q in ["🐢 海龟看见一个塑料袋,","   以为是好吃的水母…","😣 吃下去会生病,很难受。","💙 我们少用塑料,就能帮海龟!"]:
        ap(tf,q,sz=14,c=DARK)
    frame_bar(s,4.78,"我想对海龟说:____ 。","I want to tell the turtle: ____ .")
    page(s)

def s_single_use():
    cards4_slide("🥤 任务 3 · 什么是一次性塑料? Single-use plastic",
        "用一次就扔掉的塑料 — used once, then thrown away",
        [("🥤","吸管","straw"),("🛍️","塑料袋","bag"),("🍶","塑料瓶","bottle"),("🥡","塑料杯","cup")],
        frame=("____ 只用一次,就变成垃圾。","____ is used once, then becomes trash."),
        head_c=SEAGRN,card_c=SEAGRN)

def s_compare():
    s=ns(); bg(s,CREAM); hb(s,"♻️ 任务 3 · 用一次 vs 可重复 · Reusable is better",c=SEAGRN)
    kicker(s,"哪个更好? 可以用很多次的更好! Reusable is better")
    # left single-use (coral)
    box(s,0.4,1.3,4.5,3.25,fill=WHITE,line=CORAL,lw=2.5)
    pill(s,0.4,1.3,4.5,0.55,"🚫 一次性 Single-use — 用一次",CORAL,sz=14)
    su=[("🥤","吸管"),("🛍️","塑料袋"),("🥡","纸杯")]
    for k,(em,cn) in enumerate(su):
        tb(s,0.7+ k*1.45,2.1,1.3,0.7,em,sz=34,a=PP_ALIGN.CENTER)
        tb(s,0.7+ k*1.45,2.95,1.3,0.4,cn,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,0.6,3.6,4.1,0.6,"用完就扔 → 变成垃圾、伤害海洋",sz=12.5,b=True,c=CORAL,a=PP_ALIGN.CENTER)
    # right reusable (seagreen)
    box(s,5.1,1.3,4.5,3.25,fill=WHITE,line=SEAGRN,lw=2.5)
    pill(s,5.1,1.3,4.5,0.55,"✅ 可重复 Reusable — 用很多次",SEAGRN,sz=14)
    ru=[("🧴","水壶"),("👜","布袋"),("🍱","饭盒")]
    for k,(em,cn) in enumerate(ru):
        tb(s,5.4+ k*1.45,2.1,1.3,0.7,em,sz=34,a=PP_ALIGN.CENTER)
        tb(s,5.4+ k*1.45,2.95,1.3,0.4,cn,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,5.3,3.6,4.1,0.6,"用很多次 → 垃圾更少、保护海洋",sz=12.5,b=True,c=SEAGRN,a=PP_ALIGN.CENTER)
    frame_bar(s,4.78,"我选 ____ ,因为它可以用很多次。","I choose ____ because it can be reused.")
    page(s)

def s_reduce():
    cards4_slide("💪 任务 4 · 减少塑料的 4 个方法 · 4 ways",
        "海洋小卫士这样做! Ocean Guardians do this",
        [("🧴","自带水壶","bring a bottle"),("👜","自带购物袋","bring a bag"),("🚫","不用吸管","skip the straw"),("♻️","重复使用","reuse")],
        frame=("我会 ____ ,保护海洋。","I will ____ to protect the ocean."),
        head_c=DEEP,card_c=DEEP)

def s_pledge():
    s=ns(); bg(s,CREAM); hb(s,"🛡️ 任务 4 · 海洋小卫士承诺 Ocean Guardian Pledge",c=DEEP)
    kicker(s,"举起手,大声说出你的承诺! Raise your hand and say your pledge")
    box(s,1.1,1.45,7.8,1.7,fill=WHITE,line=DEEP,lw=3)
    tb(s,1.1,1.75,7.8,0.7,"「我是海洋小卫士!我会 ______ ,保护海洋。」",sz=22,b=True,c=DEEP,a=PP_ALIGN.CENTER)
    tb(s,1.1,2.55,7.8,0.4,"I'm an Ocean Guardian! I will ______ to protect the ocean.",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
    chips=[("🧴","自带水壶"),("👜","自带布袋"),("🚫","不用吸管"),("♻️","重复使用")]
    w=2.0; gap=0.2; x=(10-(w*4+gap*3))/2
    for k,(em,cn) in enumerate(chips):
        pill(s,x+k*(w+gap),3.45,w,0.7,f"{em} {cn}",SEAGRN,sz=14)
    page(s)

def s_vocab_overview():
    s=ns(); bg(s,CREAM); hb(s,"📖 我会认 · 5 个新词 5 new words")
    kicker(s,"水 · 海洋 · 塑料 · 污染 · 保护 — 今天的新词 Today's new words")
    words=[("💧","水","shuǐ",OCEAN),("🌊","海洋","hǎi yáng",SKY),("🥤","塑料","sù liào",SEAGRN),
           ("🛢️","污染","wū rǎn",CORAL),("🛡️","保护","bǎo hù",DEEP)]
    w=1.74; gap=0.16; x=(10-(w*5+gap*4))/2; y=1.45
    for k,(em,cn,py,c) in enumerate(words):
        box(s,x+k*(w+gap),y,w,2.7,fill=WHITE,line=c,lw=2.5)
        tb(s,x+k*(w+gap),y+0.2,w,0.7,em,sz=36,a=PP_ALIGN.CENTER)
        tb(s,x+k*(w+gap),y+1.05,w,0.6,cn,sz=26,b=True,c=c,a=PP_ALIGN.CENTER)
        tb(s,x+k*(w+gap),y+1.85,w,0.4,py,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
    page(s)

def s_bamboozle():
    s=ns(); bg(s,CREAM); hb(s,"🎮 复习游戏 · Bamboozle 大闯关 Team Review Game")
    kicker(s,"分组抢答,复习上午 + 今天! Review morning + today in teams")
    box(s,0.4,1.3,4.3,3.3,fill=WARM,line=SUNYEL,lw=2.5)
    tb(s,0.4,1.7,4.3,1.0,"🌊 🥤 ♻️",sz=46,a=PP_ALIGN.CENTER)
    box(s,1.0,3.0,3.1,0.75,fill=CORAL)
    tb(s,1.0,3.16,3.1,0.45,"▶ 打开 Bamboozle",sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.95,4.1,0.5,"baamboozle.com — 老师课前创建游戏",sz=11,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    box(s,4.95,1.3,4.65,2.5,fill=WHITE,line=OCEAN,lw=2.5)
    tf=tb(s,5.2,1.5,4.2,0.4,"📋 老师准备 Setup",sz=14,b=True,c=OCEAN)
    for st in ["1. 登录 baamboozle.com","2. 导入题库 day2_water_baamboozle.csv","3. 分 2–3 组,轮流抢答","4. 答对得分,复习水 + 塑料!"]:
        ap(tf,st,sz=13,c=DARK)
    pill(s,4.95,4.0,4.65,0.6,"📂 约 10 题 · 覆盖 水 / 海洋 / 塑料 / 减塑",SEAGRN,sz=13)
    page(s)

def s_learned():
    s=ns(); bg(s,CREAM); hb(s,"⭐ 我们学到 · 4 个任务完成! · 4 missions done")
    kicker(s,"海洋小卫士,你做到了! Ocean Guardian — you did it!")
    tasks=[(1,"水很重要","Water matters","💧",OCEAN),
           (2,"塑料伤害海洋","Plastic hurts the sea","🌊",CORAL),
           (3,"少用一次性","Less single-use","🥤",SEAGRN),
           (4,"我能减塑","I can reduce plastic","💪",DEEP)]
    w=2.18; gap=0.18; x=(10-(w*4+gap*3))/2
    for k,(n,cn,en,em,c) in enumerate(tasks):
        mission_card(s,x+k*(w+gap),1.45,w,2.5,n,cn,en,em,c)
    frame_bar(s,4.55,"我学到了 ____ 。下次我会 ____ 。","I learned ____ . Next time I will ____ .")
    page(s)

def s_close():
    s=ns(); bg(s,DEEP)
    box(s,0,4.55,10,1.07,fill=OCEAN,rad=False); box(s,0,4.35,10,0.25,fill=SKY,rad=False)
    tb(s,1,1.25,8,1.0,"🛡️ 我是海洋小卫士!",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.45,8,0.6,"我会保护水,减少塑料!",sz=24,b=True,c=SUNYEL,a=PP_ALIGN.CENTER)
    tb(s,1,3.25,8,0.5,"I'm an Ocean Guardian — I protect water and use less plastic!",sz=14,c=RGBColor(0xCF,0xE6,0xF0),a=PP_ALIGN.CENTER)
    tb(s,1,3.95,8,0.5,"🌊 💧 🐢 ♻️ 💙",sz=28,a=PP_ALIGN.CENTER)
    NP[0]+=1

# ================================================================ build order
def build():
    NP[0]=0
    s_cover()
    divider("复习 + 新任务","Afternoon — Review & New Mission",emoji="🌊",c=DEEP,tag="下午 AFTERNOON")
    s_recap()
    s_mission()
    s_goals()
    # 任务 1
    divider("任务 1 · 水很重要","Mission 1 — Water Matters",emoji="💧",c=OCEAN,tag="MISSION 1")
    s_water_need(); s_water_precious()
    # 任务 2
    divider("任务 2 · 海洋告急","Mission 2 — Ocean in Trouble",emoji="🌊",c=CORAL,tag="MISSION 2")
    s_plastic_to_sea(); s_animals_hurt(); s_turtle_story()
    # 任务 3
    divider("任务 3 · 一次性的麻烦","Mission 3 — Single-use Trouble",emoji="🥤",c=SEAGRN,tag="MISSION 3")
    s_single_use(); s_compare()
    # 任务 4
    divider("任务 4 · 我能做到","Mission 4 — I Can Help",emoji="💪",c=DEEP,tag="MISSION 4")
    s_reduce(); s_pledge()
    # 语言目标
    divider("语言目标 · 我会认","Language Goals — Recognize",emoji="📖",c=OCEAN,tag="LANGUAGE")
    s_vocab_overview()
    VOC=[("水","shuǐ","water","喝水","💧",OCEAN),
         ("海洋","hǎi yáng","ocean","大海","🌊",SKY),
         ("塑料","sù liào","plastic","塑料袋","🥤",SEAGRN),
         ("污染","wū rǎn","pollution","水污染","🛢️",CORAL),
         ("保护","bǎo hù","protect","保护海洋","🛡️",DEEP)]
    for i,(ch,py,en,word,em,c) in enumerate(VOC): vocab_slide(i+1,len(VOC),ch,py,en,word,em,c)
    divider("语言目标 · 我会写","Language Goals — Write",emoji="✍️",c=SEAGRN,tag="LANGUAGE")
    WR=[("水","water","4 画"),("保护","protect","保 9 画 + 护 7 画"),("塑料","plastic","塑 13 画 + 料 10 画")]
    for i,(w,en,st) in enumerate(WR): write_slide(i+1,len(WR),w,en,st)
    # game & wrap
    s_bamboozle()
    s_learned()
    s_close()
    prs.save(OUT)
    print("SAVED",OUT,len(prs.slides._sldIdLst),"slides")

if __name__=="__main__":
    build()
