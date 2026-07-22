#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零废弃 Day 5 — 环保艺术展 (Eco Art Show · 变废为宝 Trash → Treasure)
Story-driven K-5 project day, built on the 野外生存与探险 Day-3 structure.

Moved here from the 3R lesson (create_zerowaste_3R.py):
  • 「变废为宝」视频 + Project 过渡
  • 3 个分级项目: 环保拼贴画 / 瓶子变变变 / 环保海报·发明
Plus a new 环保艺术展 gallery-walk + present + celebrate closing.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Zero-Waste palette (shared with the 3R deck) ---
ECO     = RGBColor(0x2E,0x7D,0x32)
LEAF    = RGBColor(0x66,0xBB,0x6A)
DEEP    = RGBColor(0x1B,0x5E,0x20)
REDUCE  = RGBColor(0x1E,0x88,0xE5)
REUSE   = RGBColor(0xFB,0x8C,0x00)
RECYCLE = RGBColor(0x43,0xA0,0x47)
SUN     = RGBColor(0xF5,0xC2,0x42)
BROWN   = RGBColor(0x6B,0x44,0x23)
SKY     = RGBColor(0x4A,0x90,0xD9)
ALERT   = RGBColor(0xD0,0x4A,0x3C)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xF1,0xF8,0xE9)
CREAM   = RGBColor(0xFD,0xFB,0xF3)
IMGBG   = RGBColor(0xE8,0xE8,0xE8)
GOLD    = RGBColor(0xF9,0xA8,0x25)
OK      = RGBColor(0x38,0x8E,0x3C)

BASE = "/Users/Huan/0 projects/summercourse/Chinese/zero_waste零废弃"

# === Helpers (identical scaffolding to the 3R deck / 野外生存与探险 Day 3) ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi';return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name='KaiTi'
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def ib(s,l,t,w,h,lb="📷"):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def img(s,l,t,w,h,path,fallback="📷"):
    if os.path.exists(path): s.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(h))
    else: ib(s,l,t,w,h,fallback)
def hb(s,txt,c=ECO,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,0.5,1.5,9,1.2,f"{emoji} {title}",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.8,9,0.8,sub,sz=20,c=WHITE,a=PP_ALIGN.CENTER);return s
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)
def teacher_student_bar(s,t,teacher_q,student_action,color=ECO):
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=color;sf.line.width=Pt(2)
    tb(s,0.45,t+0.04,4.5,0.25,"👩‍🏫 老师问 Teacher asks:",sz=10,b=True,c=color)
    tb(s,0.45,t+0.27,4.5,0.28,teacher_q,sz=12,b=True,c=DARK)
    tb(s,5.0,t+0.04,4.6,0.25,"🧒 学生 Student does:",sz=10,b=True,c=REUSE)
    tb(s,5.0,t+0.27,4.6,0.28,student_action,sz=12,b=True,c=DARK)

# ========================================================================
#                              SLIDES
# ========================================================================
n=0

# 1. COVER
s=ns();bg(s,DEEP)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.4,8,0.5,"DAY 5",sz=18,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"🎨 环保艺术展",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"Eco Art Show  ·  变废为宝 Trash → Treasure",sz=20,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"♻️ 环保小艺术家任务  Eco Artist Mission",sz=24,b=True,c=ECO,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"看视频 · 做作品 · 说 3R · 办展览",sz=14,b=True,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"Watch · Create · Explain · Exhibit",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"零废弃 · Zero Waste Unit",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「小艺术家们! 今天我们把废物变成宝贝, 办一场环保艺术展!」\n• 收一收上次请学生带的干净废物 (瓶子/纸箱/瓶盖/旧杂志)。")

# 2. GOALS
s=ns();bg(s,CREAM);hb(s,"🎯 今天的目标  Today's Goals",ECO)
tb(s,0.4,0.85,9.2,0.4,"上完今天, 你可以…",sz=18,b=True,c=ECO,a=PP_ALIGN.CENTER)
goals=[
    ("1️⃣","复习 3R","Review Reduce · Reuse · Recycle",RECYCLE),
    ("2️⃣","用废物做一件作品","Make one thing from clean trash",REUSE),
    ("3️⃣","说出作品用了哪个 R","Say which R your work uses",REDUCE),
    ("4️⃣","在艺术展介绍你的作品","Present it at the Eco Art Show",ECO),
]
for i,(num,cn,en,cl) in enumerate(goals):
    y=1.55+i*0.8
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(y),Inches(9.2),Inches(0.68))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.55,y+0.13,0.7,0.4,num,sz=24,a=PP_ALIGN.CENTER)
    tb(s,1.30,y+0.08,4.2,0.4,cn,sz=18,b=True,c=cl)
    tb(s,5.60,y+0.14,4.0,0.4,en,sz=12,c=DARK)
n+=1;pn(s,n)
notes(s,"目标 (1 分钟):\n• 一起读 4 个目标。\n• 重点: 不只是做手工 — 要能说出用了哪个 R + 为什么。")

# 3. 3R QUICK RECAP (bridge from the 3R lesson)
s=ns();bg(s,CREAM);hb(s,"🔁 还记得 3R 吗?  Remember the 3R?",ECO)
tb(s,0.4,0.85,9.2,0.4,"今天我们主要用 Reuse — 让废物再活一次!",sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
threeR=[("⬇️","Reduce","减少","用得少一点",REDUCE),
        ("🔁","Reuse","重复使用","再用一次 ⭐今天",REUSE),
        ("♻️","Recycle","回收","变成新东西",RECYCLE)]
for i,(em,en,cn,d,cl) in enumerate(threeR):
    x=0.4+i*3.15
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.40),Inches(3.0),Inches(2.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(3)
    tb(s,x+0.1,1.60,2.8,0.9,em,sz=60,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.60,2.8,0.4,en,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.05,2.8,0.4,cn,sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,3.55,2.8,0.4,d,sz=13,c=GRAY,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.45,"3R 是哪 3 个? 今天用哪个最多?","Reduce! Reuse ⭐! Recycle! — 今天多用 Reuse!")
n+=1;pn(s,n)
notes(s,"复习 (2 分钟):\n• 快速回顾上次学的 3R。\n• 点明今天「变废为宝」主要是 Reuse (也可 Recycle 的材料再创作)。")

# 4. 变废为宝 VIDEO + PROJECT LAUNCH  (moved from the 3R lesson)
s=ns();bg(s,DEEP)
tb(s,1,0.30,8,0.6,"🎬 变废为宝  From Trash to Treasure",sz=28,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,0.90,8,0.35,"看看别人用废物做出了什么! / What can old things become?",sz=14,b=True,c=WARM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(1.35),Inches(9.0),Inches(2.5))
sh.fill.solid();sh.fill.fore_color.rgb=DARK;sh.line.color.rgb=WHITE;sh.line.width=Pt(3)
tb(s,0.5,1.90,9.0,1.2,"▶️",sz=120,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,0.5,3.35,9.0,0.40,"🔗 老师在这里插入「变废为宝」视频 / Teacher: paste video here",sz=13,b=True,c=SUN,a=PP_ALIGN.CENTER)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(4.05),Inches(9.0),Inches(1.05))
sh2.fill.solid();sh2.fill.fore_color.rgb=WARM;sh2.line.color.rgb=SUN;sh2.line.width=Pt(2)
tb(s,0.65,4.13,8.7,0.3,"👂 看的时候想:",sz=12,b=True,c=ECO)
tb(s,0.65,4.43,4.5,0.3,"1. 他们用了哪些废物材料?",sz=12,c=DARK)
tb(s,5.15,4.43,4.5,0.3,"2. 哪些材料可以 Reuse?",sz=12,c=DARK)
tb(s,0.65,4.75,9.0,0.3,"➜ 今天我们也来做! 每人一件作品, 办一场「环保艺术展」🎨",sz=13,b=True,c=ECO)
n+=1;pn(s,n)
notes(s,"视频 + 启动 (5 分钟):\n• 找 1 个 K-5 友好的「变废为宝 / trash to art / recycled crafts」视频 (2-3 分钟)。\n• 看前布置任务: 数一数用了哪些废物材料? 哪些可以 Reuse?\n• 看完引出今天的项目: 每人用废物做一件作品 → 环保艺术展。")

# 5. PROJECTS OVERVIEW (3 leveled)  (moved)
s=ns();bg(s,CREAM);hb(s,"🛠️ 动手时间! 变废为宝  Hands-On — 3 Projects",BROWN)
tb(s,0.4,0.80,9.2,0.35,"选一个适合你的项目 — 都为「环保艺术展」做准备!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
projects=[
    ("PROJECT 1","🖼️ 环保拼贴画","Eco Collage","废料贴出图案\nGlue scraps → picture",WARM,RECYCLE,"基础 K-1"),
    ("PROJECT 2","🍶 瓶子变变变","Bottle Makeover","1 个瓶子变有用\n1 bottle → useful",RGBColor(0xFF,0xE0,0xB2),REUSE,"中级 G2-3"),
    ("PROJECT 3","📢 环保海报/发明","Poster / Invention","设计 + 口号 + 说 R\nDesign + slogan",RGBColor(0xBB,0xDE,0xFB),REDUCE,"进阶 G4-5"),
]
for i,(lbl,nm,en,d,bgc,cl,lvl) in enumerate(projects):
    x=0.3+i*3.2
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.20),Inches(3.1),Inches(3.9))
    sh.fill.solid();sh.fill.fore_color.rgb=bgc;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.30,2.9,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.65,2.9,0.6,nm,sz=19,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.25,2.9,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    pill(s,x+0.75,2.65,1.6,0.35,lvl,cl,sz=10)
    ib(s,x+0.2,3.15,2.7,1.1,"📷 示范")
    ls=d.split('\n')
    tb(s,x+0.15,4.35,2.85,0.35,ls[0],sz=12,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.15,4.70,2.85,0.35,ls[1],sz=10,c=GRAY,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"项目总览 (2 分钟):\n• 3 个项目难度不同, 都低 prep, 都用「干净的废物」。\n• 可以按年级分, 也可让学生自选。\n• 全部作品用于今天的「环保艺术展」。\n• 材料: 上次请学生带的 干净塑料瓶 / 纸箱 / 瓶盖 / 旧杂志 (老师另备胶水/彩笔/剪刀)。")

# 6. PROJECT 1 — Eco Collage (Basic)  (moved)
s=ns();bg(s,CREAM);hb(s,"🖼️ Project 1: 环保拼贴画  Eco Collage  (基础 K-1)",RECYCLE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=RECYCLE;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料 (低 prep)  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.1,"🗞️ 旧杂志 / 废纸 / 纸板",sz=13,c=DARK)
ap(tf,"🔵 干净瓶盖 · 碎纸片",sz=13,c=DARK)
ap(tf,"📄 一张底纸  Base paper",sz=13,c=DARK)
ap(tf,"🖇️ 胶水  Glue",sz=13,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=REUSE;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.1,"1️⃣ 想一个图案 (地球/花/动物)",sz=13,c=DARK)
ap(tf2,"2️⃣ 撕 / 剪废料成小块",sz=13,c=DARK)
ap(tf2,"3️⃣ 拼 + 贴出图案",sz=13,c=DARK)
ap(tf2,"4️⃣ 说说你用了什么废料",sz=13,b=True,c=RECYCLE)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.35))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=RECYCLE;sh3.line.width=Pt(2)
tb(s,0.5,3.95,9,0.35,"🗣️ 展示句型  Say These:",sz=14,b=True,c=RECYCLE)
tb(s,0.5,4.35,4.5,0.35,"·  我用 ___ 做了 ___ 。",sz=14,c=DARK)
tb(s,0.5,4.70,4.5,0.35,"·  这是 Reuse! 重复使用!",sz=14,c=DARK)
tb(s,5.2,4.35,4.5,0.35,"·  I made a ___ from ___ .",sz=14,c=DARK)
tb(s,5.2,4.70,4.5,0.35,"·  It saves the Earth! 🌍",sz=14,c=DARK)
n+=1;pn(s,n)
notes(s,"Project 1 (低 prep, 15 分钟):\n• 最适合 K-1: 只需废纸/胶水, 不用剪刀也行 (撕纸)。\n• 图案给参考: 地球、树、花、小动物。\n• 完成后每人说一句「我用 ___ 做了 ___ 」。\n• 作品放进艺术展。")

# 7. PROJECT 2 — Bottle Makeover (Mid)  (moved)
s=ns();bg(s,CREAM);hb(s,"🍶 Project 2: 瓶子变变变  Bottle Makeover  (中级 G2-3)",REUSE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.4),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=REUSE;sh.line.fill.background()
tb(s,0.4,0.98,4.2,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.4,2.1,"🍶 1 个干净塑料瓶 / 罐子",sz=13,c=DARK)
ap(tf,"✂️ 剪刀 (老师帮忙剪)",sz=13,c=DARK)
ap(tf,"🖍️ 彩笔 / 贴纸 装饰",sz=13,c=DARK)
ap(tf,"🌱 泥土 + 种子 (花盆版, 可选)",sz=13,c=DARK)
ap(tf,"⚠️ 剪瓶子由老师做!",sz=12,b=True,c=ALERT)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=RECYCLE;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.2,"1️⃣ 洗干净瓶子",sz=13,c=DARK)
ap(tf2,"2️⃣ 选变成什么: 笔筒 / 花盆 / 存钱罐",sz=13,c=DARK)
ap(tf2,"3️⃣ 老师帮剪出形状",sz=13,c=DARK)
ap(tf2,"4️⃣ 装饰它 — 画画 / 贴纸",sz=13,c=DARK)
ap(tf2,"5️⃣ 完成! 旧瓶子有新工作了 🎉",sz=13,b=True,c=REUSE)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.85),Inches(9.4),Inches(1.35))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=REUSE;sh3.line.width=Pt(2)
tb(s,0.5,3.95,9,0.35,"🗣️ 展示句型  Say These:",sz=14,b=True,c=REUSE)
tb(s,0.5,4.35,4.5,0.35,"·  我把瓶子变成了 ___ 。",sz=14,c=DARK)
tb(s,0.5,4.70,4.5,0.35,"·  旧瓶子, 重复使用!",sz=14,c=DARK)
tb(s,5.2,4.35,4.5,0.35,"·  My bottle became a ___ .",sz=14,c=DARK)
tb(s,5.2,4.70,4.5,0.35,"·  That's Reuse! ♻️",sz=14,c=DARK)
n+=1;pn(s,n)
notes(s,"Project 2 (中等 prep, 20 分钟):\n• 提前请学生带 1 个洗净的塑料瓶。\n• 安全: 剪瓶口由老师用剪刀 / 美工刀完成, 边缘可贴胶带。\n• 3 个方向: 笔筒 (剪一半) / 花盆 (装土种豆) / 存钱罐 (剪小口)。\n• 花盆版可当「班级绿植」延续观察。\n• 完成后说「我把瓶子变成了 ___」。")

# 8. PROJECT 3 — Eco Poster / Invention (Advanced)  (moved)
s=ns();bg(s,CREAM);hb(s,"📢 Project 3: 环保海报 / 小发明  Poster / Invention  (进阶 G4-5)",REDUCE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.6),Inches(2.85))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=REDUCE;sh.line.width=Pt(2.5)
tb(s,0.45,1.05,4.4,0.4,"🅰️ 环保海报  Eco Poster",sz=16,b=True,c=REDUCE)
tf=tb(s,0.45,1.55,4.4,0.4,"• 画一张 Zero Waste 海报",sz=13,c=DARK)
ap(tf,"• 写一句口号 (Slogan)",sz=13,c=DARK)
ap(tf,"• 画上 3R 图标",sz=13,c=DARK)
ap(tf,"• 例: 「少一点垃圾, 多一点绿!」",sz=12,b=True,c=ECO)
ap(tf,"🧺 材料: 纸 + 彩笔 (超低 prep)",sz=12,c=GRAY)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(0.95),Inches(4.6),Inches(2.85))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=REUSE;sh2.line.width=Pt(2.5)
tb(s,5.25,1.05,4.4,0.4,"🅱️ 环保小发明  Eco Invention",sz=16,b=True,c=REUSE)
tf2=tb(s,5.25,1.55,4.4,0.4,"• 用废物做一个「有用的东西」",sz=13,c=DARK)
ap(tf2,"• 给它起名字 + 说功能",sz=13,c=DARK)
ap(tf2,"• 说它用了哪个 R",sz=13,c=DARK)
ap(tf2,"• 例: 牛奶盒 → 手机架 (Reuse)",sz=12,b=True,c=ECO)
ap(tf2,"🧺 材料: 各种干净废物 + 胶带",sz=12,c=GRAY)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.90),Inches(9.4),Inches(1.30))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=REDUCE;sh3.line.width=Pt(2)
tb(s,0.5,4.00,9,0.35,"🗣️ 展示句型  Present It:",sz=14,b=True,c=REDUCE)
tb(s,0.5,4.40,4.5,0.35,"·  我的口号是 ___ 。",sz=14,c=DARK)
tb(s,0.5,4.75,4.5,0.35,"·  我的发明用了 ___ (R)。",sz=14,c=DARK)
tb(s,5.2,4.40,4.5,0.35,"·  My slogan is ___ .",sz=14,c=DARK)
tb(s,5.2,4.75,4.5,0.35,"·  It uses Reduce / Reuse / Recycle.",sz=14,c=DARK)
n+=1;pn(s,n)
notes(s,"Project 3 (低 prep, 20-25 分钟):\n• 给 G4-5 更多思考挑战 — 二选一。\n• A 海报: 只要纸+笔; 重点是口号要短、有力。\n• B 发明: 用废物做出真正「有用」的东西, 并解释用哪个 R + 为什么。\n• 都要口头 present (1-2 句), 用词: 环保 / 减少 / 回收。\n• 成品进环保艺术展。")

# 9. THE ECO ART SHOW — gallery walk + present
s=ns();bg(s,CREAM);hb(s,"🖼️ 环保艺术展开幕!  The Eco Art Show Opens!",ECO)
tb(s,0.4,0.85,9.2,0.4,"把作品摆出来 — 一起逛展、介绍、点赞!",sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
steps=[("1️⃣","摆展台","Set up","把作品摆在桌上\n+ 写小标签",ECO),
       ("2️⃣","逛一逛","Gallery walk","安静地看每件作品\n👀",REUSE),
       ("3️⃣","做介绍","Present","轮流介绍自己的作品\n🗣️",REDUCE),
       ("4️⃣","点个赞","Vote","给喜欢的作品\n贴一颗 ⭐",GOLD)]
for i,(num,cn,en,d,cl) in enumerate(steps):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(2.20),Inches(2.55))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.55,2.0,0.5,num,sz=26,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.10,2.0,0.4,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.55,2.0,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    ls=d.split('\n')
    tf=tb(s,x+0.1,2.95,2.0,0.4,ls[0],sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    if len(ls)>1: ap(tf,ls[1],sz=14,c=DARK,a=PP_ALIGN.CENTER)
teacher_student_bar(s,4.20,"介绍一下你的作品 — 它用了哪个 R?","「这是我的 ___ 。我用 ___ 做的, 它用了 ___ (R)。」")
tb(s,0.4,4.85,9.2,0.3,"💬 This is my ___ . I made it from ___ . It uses Reduce / Reuse / Recycle.",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"环保艺术展 (10-15 分钟):\n• 学生把作品摆在桌上, 配一张小标签 (名字 + 用了哪个 R)。\n• Gallery walk: 全班安静走一圈看作品。\n• 每人用句型介绍 1-2 句 (练 Session 2 词: 环保/减少/回收)。\n• 每人得 1-2 颗 ⭐ 贴纸, 贴给喜欢的作品 (不能贴自己)。\n• 可拍照做班级「环保艺术展」留念 / 发给家长。")

# 10. CLOSING — Eco Artist badge
s=ns();bg(s,DEEP)
tb(s,1,0.6,8,0.7,"🏆 展览成功!",sz=44,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.35,8,0.5,"A Great Show!",sz=20,c=WARM,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.5),Inches(2.0),Inches(3.0),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=GOLD;sh.line.color.rgb=WHITE;sh.line.width=Pt(4)
tb(s,3.5,2.4,3.0,0.6,"🎨",sz=80,a=PP_ALIGN.CENTER)
tb(s,3.5,3.6,3.0,0.5,"环保小艺术家",sz=20,b=True,c=DEEP,a=PP_ALIGN.CENTER)
tb(s,3.5,4.1,3.0,0.4,"Eco Artist",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"记住: 少一点垃圾, 多一点绿!  Reduce · Reuse · Recycle ♻️",sz=14,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)

# === Save ===
out=os.path.join(BASE,"day5_ecoart.pptx")
prs.save(out)
print(f"Saved {out}  ({n} slides)")
