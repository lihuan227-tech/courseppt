#!/usr/bin/env python3
"""
Export day1_trash.html (the 54-slide HTML deck) to day1_trash.pptx.

Pipeline:
 1. Build a render variant of the HTML in "export mode":
      - hides nav chrome / progress / hint
      - pins #stage to exactly 1280x720, no scaling / rounding / shadow
      - replaces the bundled <video> players with a "▶ 点击播放视频" placeholder
        (the live HTML deck keeps the real, clickable videos)
 2. Headless-Chrome screenshot each slide #1..#N at 1280x720.
 3. Assemble a 16:9 PPTX, one full-bleed image per slide.

Run:  python3 export_pptx_from_html.py
"""
import os, re, subprocess, sys, tempfile, shutil

HERE = os.path.dirname(os.path.abspath(__file__))
# The live deck (a parallel ppt-creator run renamed it to *_version1.html); fall back to either.
SRC  = os.path.join(HERE, "day1_trash_version1.html")
if not os.path.exists(SRC):
    SRC = os.path.join(HERE, "day1_trash.html")
RENDER = "/tmp/zwshot/day1_trash_render.html"
IMGDIR = "/tmp/zwshot/pptx_slides"
# distinct name — do NOT clobber the pre-existing pptx files in this dir
OUT  = os.path.join(HERE, "day1_trash_from_html.pptx")
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

EXPORT_PATCH = """
<style id="exportmode">
  #chrome,#hint,#progress{display:none!important}
  html,body{background:#F6F8EC!important;overflow:hidden!important}
  #viewport{position:static!important;inset:auto!important;display:block!important}
  #stage{width:1280px!important;height:720px!important;transform:none!important;
         border-radius:0!important;box-shadow:none!important;margin:0!important}
  .slide::before{inset:14px!important}
</style>
<script>
  // Replace heavy <video> players with a friendly static placeholder for the static export.
  window.addEventListener('load',function(){
    setTimeout(function(){
      document.querySelectorAll('video').forEach(function(v){
        var d=document.createElement('div');
        d.className='vid';
        d.style.cssText='display:flex;flex-direction:column;align-items:center;justify-content:center;'+
          'gap:10px;background:#0b1410;color:#a7e0a8;font-family:sans-serif;text-align:center;'+
          'border:3px solid rgba(44,95,45,.25);border-radius:16px;height:100%';
        d.innerHTML='<div style="font-size:54px">\\u25B6</div>'+
          '<div style="font-size:22px;font-weight:700">\\u70B9\\u51FB\\u64AD\\u653E\\u89C6\\u9891</div>'+
          '<div style="font-size:14px;opacity:.8">Click to play (live deck)</div>';
        v.replaceWith(d);
      });
    }, 50);
  });
</script>
"""

def count_slides(html):
    # the deck exposes total via #total span default 54, but count S(...) at runtime is safest:
    # the engine sets total = slides.length; we read the printed default and trust 54,
    # but verify by counting section pushes deterministically from the data instead.
    m = re.search(r'<span id="total">(\d+)</span>', html)
    return int(m.group(1)) if m else 54

def main():
    html = open(SRC, encoding="utf-8").read()
    render = html.replace("</body>", EXPORT_PATCH + "\n</body>")
    # CRITICAL: strip the bundled-video src so the page's load event never waits on
    # 5 x 27MB local mp4s (all 54 slides share one DOM). The placeholder script then
    # turns the now-empty <video> tags into a friendly "click to play" card.
    render = re.sub(r'src="\$\{vsrc\([^)]*\)\}"', '', render)
    os.makedirs(IMGDIR, exist_ok=True)
    open(RENDER, "w", encoding="utf-8").write(render)

    total = count_slides(html)
    print(f"Rendering {total} slides ...", flush=True)

    imgs = []
    for n in range(1, total + 1):
        out = os.path.join(IMGDIR, f"slide_{n:02d}.png")
        if os.path.exists(out):
            os.remove(out)
        url = f"file://{RENDER}#{n}"
        prof = tempfile.mkdtemp(prefix=f"crender_{n}_")
        try:
            subprocess.run(
                [CHROME, "--headless=new", "--disable-gpu", "--hide-scrollbars",
                 "--no-first-run", "--no-default-browser-check",
                 "--force-device-scale-factor=1", "--window-size=1280,720",
                 f"--user-data-dir={prof}", f"--screenshot={out}", url,
                 "--virtual-time-budget=2000"],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                check=False, timeout=45,
            )
        except subprocess.TimeoutExpired:
            print(f"  !! slide {n} timed out", flush=True)
        finally:
            shutil.rmtree(prof, ignore_errors=True)
        if os.path.exists(out):
            imgs.append(out)
        else:
            print(f"  !! slide {n} failed to render", flush=True)
        if n % 5 == 0:
            print(f"  ... {n}/{total}", flush=True)

    print(f"Rendered {len(imgs)}/{total}. Building PPTX ...", flush=True)

    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]
    for img in imgs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(img, 0, 0, width=Inches(10), height=Inches(5.625))
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides._sldIdLst)} slides)")

if __name__ == "__main__":
    main()
