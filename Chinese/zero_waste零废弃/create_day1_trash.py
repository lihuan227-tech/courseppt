#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零 废 弃 与 可 持 续 发 展 · Day 1: 垃 圾 去 哪 儿 了?
3-session classroom deck for K-5 Chinese immersion summer camp.

Session 1 (11:00-11:45) — Discovery: Where does trash go?
Session 2 (2:00-2:45)   — Vocab + games
Session 3 (3:00-4:30)   — Project: 分类转盘 OR 环保钥匙牌
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = make_presentation()
DAY = EARTH_GREEN
n = 0


def arrow_down(s, x, y, w=0.35, h=0.45, color=DAY):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ============================================================
# 1 · COVER  (Slide 1 — 主题页)
# ============================================================
cover(prs, 1, "垃 圾 去 哪 儿 了?", "Where Does Our Trash Go?",
      "🗑️  🍌  📰  🥤",
      DAY,
      "你 扔 掉 的 垃 圾 最 后 去 了 哪 里?",
      "Where does the trash you throw away really end up?")
n += 1; pn(prs.slides[-1], n)
notes(prs.slides[-1], "📍 Day 1 Cover — 主 题 页\n👩‍🏫 老师 说: 「今天 我们 要 做 小 侦 探 — 找 出 垃 圾 的 秘 密!」\n👧 学生 做: 看 屏 幕, 一 起 读 标 题\n⏱️ 1 分钟")


# ============================================================
# 2 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1", "🔍 上 午 10:00–10:45 / 11:00–11:45  ·  垃 圾 去 哪 儿 了?",
        DAY, "🗑️"); n += 1; pn(s, n)


# ============================================================
# 3 · LEARNING GOALS
# ============================================================
s = learning_goals(prs, DAY, [
    ("1️⃣", "认 识 4 类 垃 圾 — 可 回 收/厨 余/有 害/其 他",
     "Learn 4 categories of trash", MOSS),
    ("2️⃣", "明 白 垃 圾 不 会 凭 空 消 失",
     "Trash doesn't disappear into thin air", DEEP_TEAL),
    ("3️⃣", "知 道 垃 圾 的 3 个 去 处 — 回 收/填 埋/焚 烧",
     "Know 3 places trash goes", RECYCLE_BLUE),
    ("4️⃣", "学 会 句 型: 「这 是 ___ 垃 圾」「它 可 以/不 可 以 回 收」",
     "Use the sentence frames", EARTH_GREEN),
])
n += 1; pn(s, n)


# ============================================================
# 4 · WARM UP — 透明垃圾袋 (Slide 2 — Warm Up)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗑️ 热 身  ·  看 这 个 大 袋 子!  Warm Up", DAY)

tb(s, 0.4, 0.85, 9.2, 0.28, "看 看 老 师 手 里 这 个 透 明 的 大 袋 子 — 里 面 装 了 什 么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.15, 9.2, 0.24, "Teacher holds up a clear trash bag — what's inside?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# LEFT: photo placeholder for bag
photo_slot(s, 0.40, 1.50, 4.40, 2.90,
   "📸 透 明 垃 圾 袋 + 4 件 垃 圾",
   "Clear bag with bottle/banana/paper/can", DAY)

# RIGHT: 4 items
items = [
    ("🥤", "塑 料 瓶", "Plastic bottle", RECYCLE_BLUE),
    ("🍌", "香 蕉 皮", "Banana peel", MOSS),
    ("📰", "报 纸",   "Newspaper",     EARTH_BROWN),
    ("🥫", "易 拉 罐", "Aluminum can",   FIRE_ORANGE),
]
for i, (em, cn, en, cl) in enumerate(items):
    row = i // 2; col = i % 2
    x = 5.10 + col * 2.30
    y = 1.50 + row * 1.45
    panel(s, x, y, 2.20, 1.35, cl, fill=WHITE, lw=2.5)
    tb(s, x, y+0.10, 2.20, 0.55, em, sz=32, a=PP_ALIGN.CENTER)
    tb(s, x, y+0.70, 2.20, 0.32, cn, sz=14, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, y+1.05, 2.10, 0.22, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)

# Bottom: 3 questions
qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
qbox.fill.solid(); qbox.fill.fore_color.rgb = DAY
qbox.line.color.rgb = STAR; qbox.line.width = Pt(2.5)
tb(s, 0.55, 4.62, 9.0, 0.28, "🎤 老 师 问 3 个 问 题:",
   sz=11, b=True, c=STAR, a=PP_ALIGN.LEFT)
qs = "1. 你 今 天 扔 过 垃 圾 吗?    2. 垃 圾 去 哪 儿 了?    3. 垃 圾 会 消 失 吗?"
tb(s, 0.55, 4.92, 9.0, 0.42, qs, sz=13, b=True, c=WHITE, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "📍 Slide 2 · Warm Up\n👩‍🏫 老师 做: 拿 出 透 明 垃 圾 袋 (里 面 已 装 4 件 实 物 — 真 的 干 净 物 品)\n👩‍🏫 老师 说:\n  1. 「你 今 天 扔 过 垃 圾 吗?」(举 手)\n  2. 「垃 圾 去 哪 儿 了?」(回 答)\n  3. 「垃 圾 会 消 失 吗?」(大 家 一 起 想)\n👧 学生 做: 举 手 + 大 声 回 答\n⏱️ 4-5 分钟\n💡 关 键: 不 给 答 案 — 留 悬 念!")


# ============================================================
# 5 · THINK-PAIR-SHARE (Slide 3)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🤝 想 一 想 · 说 一 说  ·  Think-Pair-Share", DAY)

# Big question banner
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.05))
qb.fill.solid(); qb.fill.fore_color.rgb = DAY
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 1.08, 9.0, 0.45, "❓ 垃 圾 最 后 去 了 哪 里?",
   sz=24, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.55, 9.0, 0.35, "Where does trash finally end up?",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)

# 3-step icons
steps = [
    ("🧠", "想 一 想", "Think · 自 己 想 30 秒", DEEP_TEAL),
    ("👥", "两 人 说", "Pair · 跟 同 桌 说 1 分 钟", MOSS),
    ("🎤", "全 班 听", "Share · 几 位 同 学 分 享", FIRE_ORANGE),
]
sw = 2.85; sgap = 0.20
stotal = 3*sw + 2*sgap; sstart = (10 - stotal)/2
for i, (em, cn, en, cl) in enumerate(steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 2.25, sw, 1.85, cl, fill=WHITE, lw=2.5)
    tb(s, x, 2.40, sw, 0.70, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 3.15, sw, 0.38, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.58, sw-0.20, 0.42, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)

# Timer card
tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.05))
tbox.fill.solid(); tbox.fill.fore_color.rgb = STAR
tbox.line.color.rgb = DAY; tbox.line.width = Pt(3)
tb(s, 0.55, 4.42, 9.0, 0.40, "⏱️ 计 时 器: 30 秒 · 60 秒 · 60 秒",
   sz=16, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.85, 9.0, 0.35, "Timer: 30s think · 60s pair · 60s share  (Total ~3 min)",
   sz=11, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 3 · Think-Pair-Share\n👩‍🏫 老师 说:\n  Step 1 「先 自 己 想 30 秒」\n  Step 2 「跟 旁 边 的 同 学 说 你 的 想 法」\n  Step 3 「请 几 位 同 学 站 起 来 分 享」\n👧 学生 做: 想 → 说 → 听\n⏱️ 3-4 分钟\n💡 老师 在 黑 板 上 列 出 学 生 答 案 — 不 评 价")


# ============================================================
# 6 · 绘本导入 (Slide 4)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "📖 绘 本 时 间  ·  《垃 圾 哪 里 去 了?》", DAY)

# LEFT: book cover + youtube link
panel(s, 0.40, 0.95, 4.40, 3.90, DAY, fill=WHITE, lw=3)
tb(s, 0.40, 1.10, 4.40, 1.00, "📚", sz=70, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.20, 4.40, 0.45, "《垃 圾 哪 里 去 了?》",
   sz=20, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.68, 4.40, 0.30, "Where Does Trash Go?",
   sz=12, c=GRAY, a=PP_ALIGN.CENTER)

# YouTube link panel
yt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.20), Inches(4.00), Inches(1.50))
yt.fill.solid(); yt.fill.fore_color.rgb = INK
yt.line.color.rgb = FIRE_ORANGE; yt.line.width = Pt(2.5)
tb(s, 0.70, 3.30, 3.80, 0.30, "▶️ YouTube 视 频",
   sz=11, b=True, c=FIRE_ORANGE, a=PP_ALIGN.LEFT)
tb(s, 0.70, 3.62, 3.80, 0.50, "youtube.com/watch?v=AKoKVQb9_80",
   sz=10, b=True, c=STAR, a=PP_ALIGN.LEFT)
yt_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(4.20), Inches(3.80), Inches(0.40))
yt_btn.fill.solid(); yt_btn.fill.fore_color.rgb = FIRE_ORANGE; yt_btn.line.fill.background()
tb(s, 0.70, 4.27, 3.80, 0.30, "👆 点 击 播 放  Click to Play",
   sz=11, b=True, c=WHITE, a=PP_ALIGN.CENTER)
yt_btn.click_action.hyperlink.address = "https://www.youtube.com/watch?v=AKoKVQb9_80"

# RIGHT: pre-watch questions
panel(s, 5.10, 0.95, 4.50, 3.90, FIRE_ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.10, 0.95, 4.50, FIRE_ORANGE, "🎬 看 之 前 — 想 一 想", sz=13)
prewatch = [
    ("👀", "你 看 到 了 什 么 垃 圾?", "What trash did you see?"),
    ("🤔", "你 觉 得 垃 圾 会 去 哪 里?", "Where do you think trash goes?"),
    ("👂", "认 真 听 + 边 看 边 想!", "Watch + listen carefully!"),
]
for i, (em, cn, en) in enumerate(prewatch):
    y = 1.65 + i * 1.00
    tb(s, 5.25, y, 0.55, 0.55, em, sz=30, a=PP_ALIGN.LEFT)
    tb(s, 5.85, y+0.08, 3.60, 0.38, cn,
       sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, 5.85, y+0.48, 3.60, 0.30, en,
       sz=9, c=GRAY, a=PP_ALIGN.LEFT)

# Bottom: instruction
inst = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.40))
inst.fill.solid(); inst.fill.fore_color.rgb = DAY; inst.line.fill.background()
tb(s, 0.55, 5.06, 9.0, 0.30, "⏱️ 5-7 分 钟 · 老师 暂 停 4 次, 问 大 家 问 题",
   sz=11, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 4 · 绘 本 导 入\n👩‍🏫 老师 做:\n  1. 点 击 视 频 链 接 (在 浏 览 器 打 开)\n  2. 先 问 2 个 「看 前」 问题\n  3. 边 看 边 暂 停 — 后 面 4 页 有 暂 停 点\n👧 学生 做: 看 视 频 + 想 答 案\n⏱️ 5-7 分 钟 (含 暂 停 讨 论)\n💡 视 频 链 接: https://www.youtube.com/watch?v=AKoKVQb9_80")


# ============================================================
# 7-10 · 绘本停顿讨论 (Slides 5-8 — 4 pause discussions)
# ============================================================
pauses = [
    ("Pause 1", "👀", "你 看 到 了 什 么 垃 圾?",
     "What trash did you see?",
     "举 出 你 看 到 的 — 大 家 一 起 数!",
     "Name what you saw — let's count together!",
     MOSS),
    ("Pause 2", "🚛", "垃 圾 车 在 做 什 么?",
     "What is the garbage truck doing?",
     "它 把 垃 圾 运 走 — 运 到 哪 里 呢?",
     "It carries trash away — but to where?",
     DEEP_TEAL),
    ("Pause 3", "❓", "为 什 么 我 们 需 要 垃 圾 车?",
     "Why do we need garbage trucks?",
     "想 一 想: 没 有 垃 圾 车, 会 怎 么 样?",
     "Think: what if there were no trucks?",
     FIRE_ORANGE),
    ("Pause 4", "💭", "垃 圾 会 消 失 吗?",
     "Will the trash disappear?",
     "「消 失」 = 没 有 了. 真 的 没 有 了 吗?",
     "'Disappear' = gone. Is it really gone?",
     RECYCLE_BLUE),
]

for label, em, cn, en, hint_cn, hint_en, cl in pauses:
    s = ns(prs); bg(s, CREAM); hb(s, f"⏸️ {label} · 暂 停 讨 论  ·  Pause & Discuss", cl)

    # Pause badge top
    pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(0.95), Inches(9.20), Inches(0.50))
    pb.fill.solid(); pb.fill.fore_color.rgb = cl; pb.line.fill.background()
    tb(s, 0.55, 1.02, 9.0, 0.38, f"⏸️  {label}  —  老 师 暂 停 视 频",
       sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)

    # LEFT: snapshot placeholder
    photo_slot(s, 0.40, 1.65, 4.30, 3.20,
       "📸 视 频 暂 停 截 图",
       "Video pause snapshot", cl)

    # RIGHT: question card
    panel(s, 5.00, 1.65, 4.60, 3.20, cl, fill=WHITE, lw=3)
    tb(s, 5.10, 1.80, 4.40, 0.80, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, 5.10, 2.70, 4.40, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, 5.10, 3.20, 4.40, 0.32, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    # Hint
    hint_bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(3.65), Inches(4.30), Inches(1.05))
    hint_bx.fill.solid(); hint_bx.fill.fore_color.rgb = WARM
    hint_bx.line.color.rgb = cl; hint_bx.line.width = Pt(1.5)
    tb(s, 5.25, 3.72, 4.10, 0.30, "💡 想 一 想:",
       sz=10, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, 5.25, 4.00, 4.10, 0.35, hint_cn, sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, 5.25, 4.38, 4.10, 0.25, hint_en, sz=9, c=GRAY, a=PP_ALIGN.LEFT)

    # Bottom: timer
    tm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.00), Inches(9.20), Inches(0.40))
    tm.fill.solid(); tm.fill.fore_color.rgb = STAR; tm.line.fill.background()
    tb(s, 0.55, 5.06, 9.0, 0.30, "⏱️ 30-45 秒 讨 论 → 继 续 播 放",
       sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
    n += 1; pn(s, n)
    notes(s, f"📍 {label} · 暂 停 讨 论\n👩‍🏫 老师 做: 暂 停 视 频, 问 大 家 问 题\n👧 学生 做: 举 手 回 答\n⏱️ 30-45 秒\n💡 不 评 价 — 让 大 家 都 说")


# ============================================================
# 11 · 垃圾不会消失 — flow chart (Slide 9)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔄 垃 圾 真 的 不 见 了 吗?  Does Trash Disappear?", DAY)

tb(s, 0.4, 0.85, 9.2, 0.32, "我 们 一 起 跟 着 垃 圾 走 一 走!  Let's follow the trash!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Flow: 家 → 垃圾桶 → 垃圾车 → ???
flow = [
    ("🏠", "家",     "Home",        DAY),
    ("🗑️", "垃 圾 桶", "Trash bin",   EARTH_BROWN),
    ("🚛", "垃 圾 车", "Garbage truck", FIRE_ORANGE),
    ("❓", "?",     "Where now?",  RECYCLE_BLUE),
]
fw = 1.95; fgap_arrow = 0.30
ftotal = 4*fw + 3*fgap_arrow
fstart = (10 - ftotal)/2
for i, (em, cn, en, cl) in enumerate(flow):
    x = fstart + i*(fw + fgap_arrow)
    panel(s, x, 1.50, fw, 2.30, cl, fill=WHITE, lw=3)
    tb(s, x, 1.65, fw, 0.85, em, sz=52, a=PP_ALIGN.CENTER)
    tb(s, x, 2.60, fw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.12, fw-0.10, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Arrow between
    if i < 3:
        arrow_x = x + fw + 0.02
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(2.55), Inches(0.26), Inches(0.30))
        a.fill.solid(); a.fill.fore_color.rgb = DAY; a.line.fill.background()

# Bottom big question
qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.10), Inches(9.20), Inches(1.30))
qbox.fill.solid(); qbox.fill.fore_color.rgb = INK
qbox.line.color.rgb = STAR; qbox.line.width = Pt(3)
tb(s, 0.55, 4.25, 9.0, 0.50, "🤔 最 后 一 站 — 到 底 在 哪 里?",
   sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.80, 9.0, 0.32, "Where is that LAST step?  (We'll find out next!)",
   sz=12, b=True, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 9 · 垃 圾 不 会 消 失 (流 程 图)\n👩‍🏫 老师 说: 「跟 着 垃 圾 走 — 家 → 桶 → 车 → ??」\n  问: 「最 后 一 步 是 哪 里?」\n👧 学生 做: 大 声 猜 → 不 给 答 案\n⏱️ 2-3 分钟\n💡 留 悬 念 — 下 一 页 才 猜")


# ============================================================
# 12 · 猜一猜 (Slide 10) — 3 options
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎯 猜 一 猜!  Guess Where!", DAY)

tb(s, 0.4, 0.85, 9.2, 0.32, "垃 圾 可 能 去 哪 里? 选 一 个!  Trash could go to ... Pick one!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 3 big option cards
options = [
    ("♻️", "回 收 中 心", "Recycling Center",
     "把 旧 的 变 成 新 的!", "Old → New!", RECYCLE_BLUE),
    ("🏔️", "填 埋 场", "Landfill",
     "埋 在 大 坑 里", "Buried in big pits", EARTH_BROWN),
    ("🔥", "焚 烧 厂", "Incinerator",
     "用 火 烧 掉", "Burned with fire", FIRE_ORANGE),
]
ow = 2.85; ogap = 0.20
ototal = 3*ow + 2*ogap; ostart = (10 - ototal)/2
for i, (em, cn, en, det_cn, det_en, cl) in enumerate(options):
    x = ostart + i*(ow + ogap)
    panel(s, x, 1.30, ow, 3.30, cl, fill=WHITE, lw=3)
    tb(s, x, 1.45, ow, 1.05, em, sz=68, a=PP_ALIGN.CENTER)
    tb(s, x, 2.55, ow, 0.50, cn, sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.12, ow-0.10, 0.32, en, sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    # Detail strip
    ds = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.20), Inches(3.55), Inches(ow-0.40), Inches(0.85))
    ds.fill.solid(); ds.fill.fore_color.rgb = WARM
    ds.line.color.rgb = cl; ds.line.width = Pt(1.5)
    tb(s, x+0.20, 3.62, ow-0.40, 0.32, det_cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x+0.20, 3.98, ow-0.40, 0.28, det_en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
    # Letter badge
    bd = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.20), Inches(1.40), Inches(0.45), Inches(0.45))
    bd.fill.solid(); bd.fill.fore_color.rgb = cl; bd.line.color.rgb = STAR; bd.line.width = Pt(2)
    tb(s, x+0.20, 1.44, 0.45, 0.40, chr(65+i), sz=16, b=True, c=WHITE, a=PP_ALIGN.CENTER)

# Bottom: vote instruction
vb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.80), Inches(9.20), Inches(0.65))
vb.fill.solid(); vb.fill.fore_color.rgb = DAY; vb.line.fill.background()
tb(s, 0.55, 4.88, 9.0, 0.32, "🙋 大 声 喊 出 你 的 选 择: A、B 还 是 C?",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22, "Shout your guess: A, B, or C?",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 10 · 猜 一 猜\n👩‍🏫 老师 做: 「现 在 大 家 猜 — A? B? 还 是 C?」\n  让 选 A 的 学 生 站 起 来 / 举 手\n  让 选 B 的 / C 的 一 一 来\n👧 学生 做: 选 + 喊 出 来\n⏱️ 2 分钟\n💡 答 案: 三 个 都 对! 真 实 世 界 三 种 都 在 用")


# ============================================================
# 13 · 三种处理方式 (Slide 11)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🌍 垃 圾 的 3 种 去 处  ·  3 Places Trash Goes", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "三 个 都 对! 真 实 世 界 里 — 这 3 种 方 法 都 在 用",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

methods = [
    ("♻️", "回 收", "Recycling",
     "瓶 子 → 新 瓶 子",
     "Bottle → new bottle",
     "📸 工 人 在 分 拣 塑 料 瓶",
     "Workers sorting plastic", RECYCLE_BLUE),
    ("🏔️", "填 埋", "Landfill",
     "挖 大 坑 — 埋 起 来",
     "Big pits — buried",
     "📸 山 一 样 的 垃 圾 堆",
     "Mountain of trash", EARTH_BROWN),
    ("🔥", "焚 烧", "Incineration",
     "用 火 烧 — 变 成 灰",
     "Burned → ash",
     "📸 焚 烧 厂 + 烟 囱",
     "Incinerator chimney", FIRE_ORANGE),
]
mw = 2.90; mgap = 0.15
mtotal = 3*mw + 2*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, line_cn, line_en, photo_cn, photo_en, cl) in enumerate(methods):
    x = mstart + i*(mw + mgap)
    # Card
    panel(s, x, 1.25, mw, 4.05, cl, fill=WHITE, lw=3)
    # Header strip
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.25), Inches(mw), Inches(0.85))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.32, mw, 0.40, em + "  " + cn, sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 1.72, mw-0.10, 0.30, en, sz=10, c=WARM, a=PP_ALIGN.CENTER)
    # Photo placeholder
    photo_slot(s, x+0.15, 2.20, mw-0.30, 1.75, photo_cn, photo_en, cl)
    # One-line description
    tb(s, x+0.10, 4.05, mw-0.20, 0.35, line_cn, sz=13, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 4.40, mw-0.20, 0.28, line_en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
    # Bottom corner indicator
    tb(s, x+0.10, 4.78, mw-0.20, 0.30,
       "💚 最 好" if cl == RECYCLE_BLUE else ("⚠️ 不 够 好" if cl == EARTH_BROWN else "⚡ 有 污 染"),
       sz=11, b=True, c=cl, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 11 · 3 种 去 处\n👩‍🏫 老师 说:\n  「♻️ 回收 是 最 好 的!」「🏔️ 填埋 占 地 + 慢」「🔥 焚烧 快, 但 有 烟」\n  每 种 只 讲 一 句\n👧 学生 做: 看 图 + 听\n⏱️ 3 分钟")


# ============================================================
# 14 · 4 个 垃 圾 桶 概 览 (NEW — intro to 4-bin system)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗑️ 垃 圾 分 4 类  ·  4 种 垃 圾 桶", DAY)

tb(s, 0.4, 0.85, 9.2, 0.32, "你 知 道 吗? 垃 圾 要 分 成 4 大 类 — 每 一 类 有 自 己 的 桶!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

bins = [
    ("🟦", "可 回 收 物", "kě huí shōu wù", "Recyclable", RECYCLE_BLUE),
    ("🟩", "厨 余 垃 圾", "chú yú lā jī",   "Kitchen / Food", MOSS),
    ("🟥", "有 害 垃 圾", "yǒu hài lā jī",  "Hazardous",     RGBColor(0xC8,0x25,0x3E)),
    ("⬜", "其 他 垃 圾", "qí tā lā jī",    "Other",         RGBColor(0x60,0x60,0x60)),
]
bw = 2.20; bgap = 0.15
btotal = 4*bw + 3*bgap; bstart = (10 - btotal)/2
for i, (em, cn, py, en, cl) in enumerate(bins):
    x = bstart + i*(bw + bgap)
    # Card with colored top strip
    panel(s, x, 1.30, bw, 3.85, cl, fill=WHITE, lw=3)
    # Top color strip
    ts = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.30), Inches(bw), Inches(0.55))
    ts.fill.solid(); ts.fill.fore_color.rgb = cl; ts.line.fill.background()
    tb(s, x, 1.38, bw, 0.40, f"#{i+1}", sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
    # Bin emoji
    tb(s, x, 1.95, bw, 1.20, em, sz=72, a=PP_ALIGN.CENTER)
    # Name
    tb(s, x, 3.20, bw, 0.45, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.65, bw-0.10, 0.28, py, sz=10, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.92, bw-0.10, 0.28, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
    # Reminder strip at bottom
    rm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.10), Inches(4.30), Inches(bw-0.20), Inches(0.70))
    rm.fill.solid(); rm.fill.fore_color.rgb = WARM
    rm.line.color.rgb = cl; rm.line.width = Pt(1.5)
    samples = ["瓶 子 · 纸 · 罐 子", "果 皮 · 剩 饭 · 茶 叶", "电 池 · 灯 泡 · 药", "卫 生 纸 · 旧 牙 刷 · 碎 碗"][i]
    tb(s, x+0.10, 4.38, bw-0.20, 0.55, samples, sz=9, b=True, c=cl, a=PP_ALIGN.CENTER)

# Bottom: next step
nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.25), Inches(9.20), Inches(0.30))
nb.fill.solid(); nb.fill.fore_color.rgb = STAR; nb.line.fill.background()
tb(s, 0.55, 5.27, 9.0, 0.25, "👉 接 下 来: 我 们 一 个 一 个 来 认 识 这 4 类!",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 4 类 垃 圾 概 览\n👩‍🏫 老师 说: 「中 国 现 在 把 垃 圾 分 4 类 — 蓝/绿/红/灰 4 个 桶」\n  让 学 生 大 声 读 一 遍 4 个 名 字\n👧 学生 做: 跟 读\n⏱️ 2-3 分钟\n💡 强 调 颜 色 — 蓝 = 可 回 收, 绿 = 厨 余, 红 = 有 害, 灰 = 其 他")


# ============================================================
# 15-18 · 4 类 垃 圾 — 一 类 一 页 (one slide per category)
# ============================================================
categories = [
    {
        "color": RECYCLE_BLUE,
        "emoji": "🟦",
        "name_cn": "可 回 收 物",
        "name_py": "kě huí shōu wù",
        "name_en": "Recyclable",
        "tag_cn": "可 以 变 成 新 东 西",
        "tag_en": "Can be made into new things",
        "items": [
            ("🥤", "塑 料 瓶"),
            ("📰", "报 纸"),
            ("📦", "纸 箱"),
            ("🥫", "易 拉 罐"),
            ("🍾", "玻 璃 瓶"),
        ],
        "say_cn": "「这 些 都 可 以 回 收 — 变 成 新 的 东 西!」",
    },
    {
        "color": MOSS,
        "emoji": "🟩",
        "name_cn": "厨 余 垃 圾",
        "name_py": "chú yú lā jī",
        "name_en": "Kitchen / Food Waste",
        "tag_cn": "厨 房 里 吃 剩 下 的",
        "tag_en": "Leftovers from the kitchen",
        "items": [
            ("🍌", "香 蕉 皮"),
            ("🍎", "苹 果 核"),
            ("🍚", "剩 饭"),
            ("🥬", "烂 菜 叶"),
            ("🍵", "茶 叶 渣"),
        ],
        "say_cn": "「厨 房 里 吃 剩 下 的 都 是 厨 余 垃 圾!」",
    },
    {
        "color": RGBColor(0xC8,0x25,0x3E),
        "emoji": "🟥",
        "name_cn": "有 害 垃 圾",
        "name_py": "yǒu hài lā jī",
        "name_en": "Hazardous Waste",
        "tag_cn": "对 人 + 地 球 有 害",
        "tag_en": "Dangerous to people & Earth",
        "items": [
            ("🔋", "废 电 池"),
            ("💡", "废 灯 泡"),
            ("💊", "过 期 药 品"),
            ("🌡️", "水 银 温 度 计"),
            ("🧴", "化 妆 品"),
        ],
        "say_cn": "「这 些 一 定 要 小 心! 不 能 乱 扔!」",
    },
    {
        "color": RGBColor(0x60,0x60,0x60),
        "emoji": "⬜",
        "name_cn": "其 他 垃 圾",
        "name_py": "qí tā lā jī",
        "name_en": "Other Waste",
        "tag_cn": "前 3 类 都 不 是, 就 是 其 他",
        "tag_en": "Not the other 3 — it's 'Other'",
        "items": [
            ("🧻", "用 过 的 卫 生 纸"),
            ("🪥", "旧 牙 刷"),
            ("🍽️", "破 碗 / 陶 瓷"),
            ("👟", "旧 鞋 子"),
            ("🦴", "大 骨 头"),
        ],
        "say_cn": "「不 知 道 是 哪 一 类? 就 扔 进 「其 他」 桶!」",
    },
]

for cat in categories:
    s = ns(prs); bg(s, CREAM)
    hb(s, f"{cat['emoji']}  {cat['name_cn']}  ·  {cat['name_en']}", cat["color"])

    # Tag line
    tb(s, 0.4, 0.85, 9.2, 0.32, f"💡 {cat['tag_cn']}",
       sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.4, 1.18, 9.2, 0.26, cat["tag_en"],
       sz=10, c=GRAY, a=PP_ALIGN.CENTER)

    # LEFT: big bin card + name + pinyin
    panel(s, 0.40, 1.55, 3.50, 3.60, cat["color"], fill=WARM, lw=3)
    tb(s, 0.40, 1.75, 3.50, 1.30, cat["emoji"], sz=110, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 3.10, 3.50, 0.55, cat["name_cn"], sz=26, b=True, c=cat["color"], a=PP_ALIGN.CENTER)
    tb(s, 0.40, 3.65, 3.50, 0.30, cat["name_py"], sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, 0.40, 3.98, 3.50, 0.30, cat["name_en"], sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    # Teacher line
    tb(s, 0.50, 4.40, 3.30, 0.30, "👩‍🏫 老 师 说:",
       sz=10, b=True, c=cat["color"], a=PP_ALIGN.LEFT)
    tb(s, 0.50, 4.70, 3.30, 0.40, cat["say_cn"],
       sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)

    # RIGHT: 5 example items
    panel(s, 4.10, 1.55, 5.50, 3.60, cat["color"], fill=WHITE, lw=3)
    panel_head(s, 4.10, 1.55, 5.50, cat["color"], "✦ 都 有 哪 些 东 西?  Examples", sz=12)
    # 5 items in a 2-row layout (3 + 2)
    for j, (em, cn) in enumerate(cat["items"]):
        row = j // 3; col = j % 3
        cell_w = 1.65; cell_gap = 0.05
        # For row 0 (3 items): center starting at x=4.20
        # For row 1 (2 items): center 2 items
        if row == 0:
            x = 4.20 + col * (cell_w + cell_gap)
        else:
            x = 4.20 + (cell_w + cell_gap) * 0.5 + col * (cell_w + cell_gap)
        y = 2.15 + row * 1.40
        cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(cell_w), Inches(1.30))
        cd.fill.solid(); cd.fill.fore_color.rgb = WARM
        cd.line.color.rgb = cat["color"]; cd.line.width = Pt(2)
        tb(s, x, y+0.15, cell_w, 0.55, em, sz=30, a=PP_ALIGN.CENTER)
        tb(s, x+0.05, y+0.78, cell_w-0.10, 0.45, cn, sz=11, b=True, c=cat["color"], a=PP_ALIGN.CENTER)

    n += 1; pn(s, n)
    notes(s, f"📍 {cat['name_cn']}\n👩‍🏫 老师 说: 「{cat['say_cn']}」\n  指 着 5 张 卡 一 个 一 个 念\n  让 学 生 跟 读 一 遍\n👧 学生 做: 看 + 跟 读\n⏱️ 2-3 分钟")


# ============================================================
# 19 · 小 组 分 类 游 戏 (updated — uses 4-bin system)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🏆 小 组 游 戏  ·  4 类 垃 圾 分 类 大 挑 战!", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "桌 子 上 有 4 个 桶 — 把 8 张 卡 片 放 到 对 的 桶 里!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 sorting zones (matches 4-bin system)
zones = [
    ("🟦", "可 回 收 物", "Recyclable", RECYCLE_BLUE),
    ("🟩", "厨 余 垃 圾", "Food",       MOSS),
    ("🟥", "有 害 垃 圾", "Hazardous",  RGBColor(0xC8,0x25,0x3E)),
    ("⬜", "其 他 垃 圾", "Other",      RGBColor(0x60,0x60,0x60)),
]
zw = 2.20; zgap = 0.15
ztotal = 4*zw + 3*zgap; zstart = (10 - ztotal)/2
for i, (em, cn, en, cl) in enumerate(zones):
    x = zstart + i*(zw + zgap)
    panel(s, x, 1.30, zw, 1.55, cl, fill=WHITE, lw=2.5)
    tb(s, x, 1.40, zw, 0.55, em, sz=34, a=PP_ALIGN.CENTER)
    tb(s, x, 2.02, zw, 0.42, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.48, zw-0.10, 0.28, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# Card row — 8 cards spanning the 4 categories
tb(s, 0.4, 3.00, 9.2, 0.32, "🎴 卡 片  (每 组 一 套 · 8 张):",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
cards = [
    ("🥤", "塑 料 瓶"),
    ("📦", "纸 箱"),
    ("🍌", "香 蕉 皮"),
    ("🍚", "剩 饭"),
    ("🔋", "废 电 池"),
    ("💊", "过 期 药"),
    ("🧻", "用 过 卫 生 纸"),
    ("🪥", "旧 牙 刷"),
]
cw = 1.05; cgap = 0.08
ctotal = 8*cw + 7*cgap; cstart = (10 - ctotal)/2
for i, (em, cn) in enumerate(cards):
    x = cstart + i*(cw + cgap)
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.40), Inches(cw), Inches(1.20))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = DAY; cd.line.width = Pt(1.5)
    tb(s, x, 3.50, cw, 0.55, em, sz=24, a=PP_ALIGN.CENTER)
    tb(s, x+0.02, 4.05, cw-0.04, 0.40, cn, sz=9, b=True, c=DAY, a=PP_ALIGN.CENTER)

# Bottom: instruction
inst = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.80), Inches(9.20), Inches(0.65))
inst.fill.solid(); inst.fill.fore_color.rgb = DAY
inst.line.color.rgb = STAR; inst.line.width = Pt(2.5)
tb(s, 0.55, 4.88, 9.0, 0.32, "👥 全 班 分 4-5 组 · 一 起 把 卡 片 放 到 对 的 桶 里!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.22, 9.0, 0.22, "Split into 4-5 groups · sort each card into the right bin",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 小 组 分 类 挑 战 (4 类)\n👩‍🏫 老师 准 备: 提 前 打 印 4-5 套 卡 片 (8 张/套) + 4 个 桶 (颜 色 不 同)\n👩‍🏫 老师 说: 「全 组 一 起 商 量 — 这 张 放 哪 个 桶?」\n👧 学生 做: 小 组 合 作 — 不 要 跑\n⏱️ 5-6 分钟\n💡 老 师 走 动 — 鼓 励 用 中 文: 「这 是 ___ 垃 圾」")


# ============================================================
# 15 · 竞赛环节 + 排名 (Slide 13)
# ============================================================
s = ns(prs); bg(s, INK); hb(s, "🏆 公 布 答 案 + 比 一 比!  Reveal & Rank", FIRE_ORANGE)

tb(s, 0.4, 0.85, 9.2, 0.30, "老 师 念 答 案 — 每 组 自 己 核 对, 数 一 数 你 们 对 了 几 张!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.15, 9.2, 0.24, "Teacher reads answers · groups check + count their correct cards",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)

# Answer key strip
ak = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.50), Inches(9.20), Inches(0.85))
ak.fill.solid(); ak.fill.fore_color.rgb = WHITE
ak.line.color.rgb = STAR; ak.line.width = Pt(2.5)
tb(s, 0.55, 1.58, 9.0, 0.28, "✅ 答 案  Answer Key:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
key1 = "🥤塑料瓶 → 可回收   📦纸箱 → 可回收   🍌香蕉皮 → 厨余   🍚剩饭 → 厨余"
key2 = "🔋废电池 → 有害   💊过期药 → 有害   🧻用过卫生纸 → 其他   🪥旧牙刷 → 其他"
tb(s, 0.55, 1.85, 9.0, 0.30, key1, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 2.13, 9.0, 0.30, key2, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)

# 3 podium cards
medals = [
    ("🥇", "第 一 名", "1st Place", GOLD_MEDAL),
    ("🥈", "第 二 名", "2nd Place", SILVER_MEDAL),
    ("🥉", "第 三 名", "3rd Place", BRONZE_MEDAL),
]
mw = 2.85; mgap = 0.20
mtotal = 3*mw + 2*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, cl) in enumerate(medals):
    x = mstart + i*(mw + mgap)
    panel(s, x, 2.55, mw, 1.95, cl, fill=WHITE, lw=3)
    tb(s, x, 2.70, mw, 0.85, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 3.55, mw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 4.05, mw-0.20, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)

# Celebration banner
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.70), Inches(9.20), Inches(0.75))
cb.fill.solid(); cb.fill.fore_color.rgb = STAR; cb.line.fill.background()
tb(s, 0.55, 4.80, 9.0, 0.35, "🎉 🎊  全 班 鼓 掌!  Big round of applause!  🎉 🎊",
   sz=15, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.20, 9.0, 0.22, "每 个 组 都 是 「环 保 小 队」!  Every team is an Eco-Squad!",
   sz=10, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 13 · 公 布 答 案\n👩‍🏫 老师 做:\n  1. 一 张 一 张 念 答 案\n  2. 每 组 数 自 己 对 了 几 张\n  3. 排 名 第 一/二/三\n  4. 全 班 鼓 掌\n👧 学生 做: 核 对 + 庆 祝\n⏱️ 3-4 分钟\n💡 重 点 不 是 输 赢 — 是 让 大 家 都 觉 得 自 己 是 环 保 小 队!")


# ============================================================
# 16 · EXIT TICKET (Slide 14)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎫 出 门 票  ·  Exit Ticket — 用 句 型 说 一 个!", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "请 5-8 位 同 学 — 拿 一 件 东 西, 用 句 型 说 一 句!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Big sentence frames
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(1.30), Inches(8.80), Inches(1.30))
fr1.fill.solid(); fr1.fill.fore_color.rgb = DAY
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.42, 8.50, 0.35, "1️⃣  句 型 一",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.78, 8.50, 0.55, "「这 是 _______ 垃 圾.」",
   sz=26, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.32, 8.50, 0.25, '"This is ___ trash."',
   sz=11, c=WARM, a=PP_ALIGN.CENTER)

fr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(2.80), Inches(8.80), Inches(1.30))
fr2.fill.solid(); fr2.fill.fore_color.rgb = MOSS
fr2.line.color.rgb = STAR; fr2.line.width = Pt(3)
tb(s, 0.75, 2.92, 8.50, 0.35, "2️⃣  句 型 二",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.28, 8.50, 0.55, "「它 可 以 / 不 可 以 回 收.」",
   sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 3.82, 8.50, 0.25, '"It can / cannot be recycled."',
   sz=11, c=WARM, a=PP_ALIGN.CENTER)

# Example row
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.10))
ex.fill.solid(); ex.fill.fore_color.rgb = WARM
ex.line.color.rgb = DAY; ex.line.width = Pt(2)
tb(s, 0.55, 4.40, 9.00, 0.30, "💡 举 个 例 子  Example:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.70, 9.00, 0.32, "🥤 「这 是 塑 料 垃 圾. 它 可 以 回 收.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.05, 9.00, 0.30, "🍌 「这 是 食 物 垃 圾. 它 不 可 以 回 收.」",
   sz=14, b=True, c=DARK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "📍 Slide 14 · Exit Ticket\n👩‍🏫 老师 做: 拿 起 一 件 东 西 (瓶子/纸/...) → 点 一 位 学 生\n  学 生 用 2 个 句 型 说 一 遍\n👧 学生 做: 看 屏 幕 + 大 声 说\n⏱️ 5-6 分钟 (5-8 位 学 生)\n💡 老 师 给 「拍 拍 手」 鼓 励")


# ============================================================
# 17 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2", "📚 下 午 2:00–2:45  ·  词 汇 + 游 戏 复 习",
        DAY, "🔤"); n += 1; pn(s, n)


# ============================================================
# 18-22 · 我 会 认 · 5 个 词 (one per slide)
# ============================================================
recog_words = [
    ("📜", "纸",    "zhǐ",   "paper",
     "我 用 纸 写 字.",
     "I use paper to write.",
     "📸 一 张 白 纸 / 报 纸 / 课 本", EARTH_BROWN),
    ("🗑️", "垃 圾", "lā jī", "trash",
     "请 把 垃 圾 扔 进 桶 里.",
     "Please throw the trash in the bin.",
     "📸 一 个 垃 圾 桶 + 垃 圾", FIRE_ORANGE),
    ("♻️", "回 收", "huí shōu", "recycle",
     "塑 料 瓶 可 以 回 收.",
     "Plastic bottles can be recycled.",
     "📸 回 收 标 志 + 蓝 色 回 收 桶", RECYCLE_BLUE),
    ("🥤", "塑 料", "sù liào", "plastic",
     "这 个 瓶 子 是 塑 料 做 的.",
     "This bottle is made of plastic.",
     "📸 一 堆 塑 料 瓶 / 塑 料 杯", MOSS),
    ("🍼", "瓶 子", "píng zi", "bottle",
     "我 喝 水 用 瓶 子.",
     "I use a bottle to drink water.",
     "📸 各 种 瓶 子 (水 瓶/牛 奶 瓶)", DEEP_TEAL),
]

for em, cn, py, en, ex_cn, ex_en, hint, cl in recog_words:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)
    notes(s, f"📍 我 会 认 · {cn}\n👩‍🏫 老师 说: 「跟 我 读 — {cn}!」 (3 遍)\n  让 学 生 用 句 子 造 一 个: 例 句 「{ex_cn}」\n👧 学生 做: 跟 读 + 造 句\n⏱️ 2-3 分钟")


# ============================================================
# 19 · 词汇图片配对 (matching game)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔗 配 对 游 戏  ·  Match the Words", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "把 字 跟 图 连 起 来!  Match each word with its picture!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT column: words; RIGHT column: pictures (random order)
left_words = [
    ("纸",    EARTH_BROWN),
    ("垃 圾", FIRE_ORANGE),
    ("回 收", RECYCLE_BLUE),
    ("塑 料", MOSS),
    ("瓶 子", DEEP_TEAL),
]
right_pics = [
    ("🥤", "C"),
    ("📜", "E"),
    ("🗑️", "A"),
    ("🍼", "D"),
    ("♻️", "B"),
]

# 5 word boxes left side
for i, (w, cl) in enumerate(left_words):
    y = 1.30 + i * 0.78
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(y), Inches(3.50), Inches(0.65))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = cl; box.line.width = Pt(2.5)
    tb(s, 0.60, y+0.10, 0.50, 0.45, str(i+1), sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, 1.15, y+0.10, 2.85, 0.45, w, sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)

# 5 picture boxes right side (shuffled order)
for i, (em, letter) in enumerate(right_pics):
    y = 1.30 + i * 0.78
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.90), Inches(y), Inches(3.50), Inches(0.65))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DAY; box.line.width = Pt(2.5)
    tb(s, 5.90, y+0.10, 0.50, 0.45, letter, sz=18, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, 6.45, y+0.05, 2.85, 0.55, em, sz=28, a=PP_ALIGN.LEFT)

# Connecting hint
tb(s, 4.15, 3.10, 1.65, 0.40, "🤝 配 对",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
arr = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(4.15), Inches(3.50), Inches(1.65), Inches(0.40))
arr.fill.solid(); arr.fill.fore_color.rgb = STAR; arr.line.fill.background()

# Bottom: answer key (small)
ak = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(5.20), Inches(9.20), Inches(0.30))
ak.fill.solid(); ak.fill.fore_color.rgb = DAY; ak.line.fill.background()
tb(s, 0.55, 5.22, 9.0, 0.25, "🔑 答 案 (老师 用): 1-E  2-A  3-B  4-C  5-D",
   sz=10, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 19 · 配 对 游 戏\n👩‍🏫 老师 说: 「左 边 5 个 词, 右 边 5 张 图 — 配 一 配!」\n  可 以 全 班 一 起 喊 答 案, 或 找 学 生 上 来 连\n👧 学生 做: 看 + 喊 答 案 (1-E / 2-A...)\n⏱️ 3-4 分钟")


# ============================================================
# (BINGO slide removed per request)
# ============================================================

# ============================================================
# 拍 词 卡 游 戏 (slap cards)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "👋 拍 词 卡!  ·  Slap the Word!", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "两 人 一 组 — 老 师 念 词, 谁 先 拍 到 谁 赢!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# Cards layout (6 big cards)
slap_cards = [
    ("📜", "纸",    EARTH_BROWN),
    ("🗑️", "垃 圾", FIRE_ORANGE),
    ("♻️", "回 收", RECYCLE_BLUE),
    ("🥤", "塑 料", MOSS),
    ("🍼", "瓶 子", DEEP_TEAL),
    ("👋", "拍!",   STAR),
]
scw = 2.85; scgap = 0.15
sctotal = 3*scw + 2*scgap; scstart = (10 - sctotal)/2
for i, (em, cn, cl) in enumerate(slap_cards):
    row = i // 3; col = i % 3
    x = scstart + col*(scw + scgap)
    y = 1.30 + row*1.40
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(scw), Inches(1.20))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = cl; cd.line.width = Pt(3)
    tb(s, x+0.15, y+0.20, 0.90, 0.85, em, sz=44, a=PP_ALIGN.LEFT)
    tb(s, x+1.10, y+0.30, scw-1.20, 0.55, cn, sz=22, b=True, c=cl, a=PP_ALIGN.LEFT)

# Bottom: how to play
how_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.40), Inches(9.20), Inches(1.00))
how_box.fill.solid(); how_box.fill.fore_color.rgb = DAY
how_box.line.color.rgb = STAR; how_box.line.width = Pt(2.5)
tb(s, 0.55, 4.50, 9.00, 0.30, "🎯 怎 么 玩  How to Play:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.80, 9.00, 0.30, "• 两 人 一 组, 桌 上 摆 5 张 词 卡   • 老师 念 词   • 谁 先 拍 中 谁 赢!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.13, 9.00, 0.25, "Pairs · 5 cards on table · teacher calls word · first slap wins",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "📍 Slide 21 · 拍 词 卡\n👩‍🏫 老师 准 备: 提 前 打 印 词 卡 (5 张/对)\n  分 配 同 学 两 人 一 组\n👩‍🏫 老师 说: 「老 师 念 词 — 谁 先 拍 谁 赢!」\n👧 学生 做: 听 + 拍 (轻 拍! 不 是 打)\n⏱️ 4-5 分钟")


# ============================================================
# 22 · 句型练习 + Pair Share
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗣️ 句 型 练 习  ·  Sentence Practice", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "和 同 桌 一 起 — 一 人 拿 一 件 东 西, 用 句 型 说 给 对 方 听!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 2 sentence frames stacked
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(1.20), Inches(8.80), Inches(1.05))
fr1.fill.solid(); fr1.fill.fore_color.rgb = DAY
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.30, 8.50, 0.30, "1️⃣  这 是 ___ 垃 圾.",
   sz=20, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.65, 8.50, 0.35, "→  这 是 (纸 / 塑 料 / 食 物 / 金 属) 垃 圾.",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.75, 2.02, 8.50, 0.20, "This is ___ trash.",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)

fr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(2.45), Inches(8.80), Inches(1.05))
fr2.fill.solid(); fr2.fill.fore_color.rgb = MOSS
fr2.line.color.rgb = STAR; fr2.line.width = Pt(3)
tb(s, 0.75, 2.55, 8.50, 0.30, "2️⃣  它 可 以 / 不 可 以 回 收.",
   sz=20, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 2.90, 8.50, 0.35, "→  它 (可 以 / 不 可 以) 回 收.",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.27, 8.50, 0.20, "It (can / cannot) be recycled.",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)

# Pair Share callout
ps = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.75), Inches(9.20), Inches(1.65))
ps.fill.solid(); ps.fill.fore_color.rgb = WARM
ps.line.color.rgb = DAY; ps.line.width = Pt(2.5)
tb(s, 0.55, 3.85, 9.00, 0.32, "👥 Pair Share  ·  和 同 桌 一 起 练!",
   sz=14, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.22, 9.00, 0.30, "🔹 A 同 学: 「这 是 塑 料 垃 圾.」",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.52, 9.00, 0.30, "🔹 B 同 学: 「它 可 以 回 收!」",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.85, 9.00, 0.30, "🔄 换 一 件 东 西 — 互 换 角 色, 再 来 一 次!",
   sz=12, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.18, 9.00, 0.22, "Swap the object · switch roles · repeat!",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "📍 Slide 22 · 句 型 + Pair Share\n👩‍🏫 老师 说: 「两 人 一 组 — 一 人 拿 东 西 说 句 1, 另 一 个 说 句 2」\n  老师 走 动 监 听\n👧 学生 做: 配 对 + 说 + 换\n⏱️ 5-6 分钟")


# ============================================================
# 我 会 写 · 3 个 词 (one per slide)
# ============================================================
write_words = [
    ("纸", "paper", EARTH_BROWN, [
        ("纸", "zhǐ", "7 笔 / 7 strokes", "纟 (绞 丝 旁) + 氏 — 跟 「丝」 有 关"),
    ]),
    ("垃 圾", "trash", FIRE_ORANGE, [
        ("垃", "lā", "8 笔 / 8 strokes", "土 字 旁 + 立 — 跟 土 / 地 有 关"),
        ("圾", "jī", "6 笔 / 6 strokes", "土 字 旁 + 及 — 也 跟 土 有 关"),
    ]),
    ("回 收", "recycle", RECYCLE_BLUE, [
        ("回", "huí", "6 笔 / 6 strokes", "口 里 面 还 有 一 个 口 — 「回 来」 的 意 思"),
        ("收", "shōu", "6 笔 / 6 strokes", "攵 (反 文 旁) — 跟 「做 事」 有 关"),
    ]),
]
for cn_phrase, en, cl, chars in write_words:
    s = vocab_write(prs, cl, cn_phrase, en, chars)
    n += 1; pn(s, n)
    notes(s, f"📍 我 会 写 · {cn_phrase}\n👩‍🏫 老师 做: 黑 板 / 屏 幕 一 笔 一 笔 示 范\n  让 学 生 在 田 字 格 里 写 3 遍\n👧 学生 做: 看 → 描 → 写\n⏱️ 3-4 分 钟 (一 个 字/词)\n💡 K-2 可 以 描 红, 3-5 自 己 写")


# ============================================================
# 24 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3", "🛠️ 下 午 3:00–4:30  ·  我 是 垃 圾 分 类 专 家!",
        DAY, "♻️"); n += 1; pn(s, n)


# ============================================================
# 25 · PROJECT MENU — choose 1 or 2
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🛠️ 今 天 做 什 么?  ·  Today's Projects", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "选 一 个 或 两 个 都 做! 一 起 来 创 造!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 2 project cards
projects = [
    ("Project 1", "🎡", "垃 圾 分 类 转 盘", "Sorting Spinner",
     "🍽️ 纸 盘 · 📌 Brass fastener · 🖍️ 彩 笔",
     "Paper plate · brass fastener · markers", MOSS),
    ("Project 2", "🔑", "环 保 钥 匙 牌", "Shrinky Dink Keychain",
     "✨ Shrinky Dink · 🖍️ 彩 笔 · 🔑 钥 匙 环 · 🔥 烤 箱",
     "Shrinky Dink · markers · keyring · oven", RECYCLE_BLUE),
]
pw = 4.40; pgap = 0.20
ptotal = 2*pw + pgap; pstart = (10 - ptotal)/2
for i, (label, em, cn, en, mat_cn, mat_en, cl) in enumerate(projects):
    x = pstart + i*(pw + pgap)
    panel(s, x, 1.30, pw, 3.95, cl, fill=WHITE, lw=3)
    # Header strip
    hd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.30), Inches(pw), Inches(0.55))
    hd.fill.solid(); hd.fill.fore_color.rgb = cl; hd.line.fill.background()
    tb(s, x, 1.40, pw, 0.40, label, sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
    # Big emoji
    tb(s, x, 1.95, pw, 1.25, em, sz=88, a=PP_ALIGN.CENTER)
    # Title
    tb(s, x, 3.30, pw, 0.50, cn, sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.82, pw-0.10, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    # Materials box
    mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.20), Inches(4.20), Inches(pw-0.40), Inches(0.95))
    mb.fill.solid(); mb.fill.fore_color.rgb = WARM
    mb.line.color.rgb = cl; mb.line.width = Pt(1.5)
    tb(s, x+0.30, 4.25, pw-0.60, 0.28, "🎒 材 料:",
       sz=10, b=True, c=cl, a=PP_ALIGN.LEFT)
    tb(s, x+0.30, 4.50, pw-0.60, 0.32, mat_cn, sz=10, b=True, c=DARK, a=PP_ALIGN.LEFT)
    tb(s, x+0.30, 4.83, pw-0.60, 0.28, mat_en, sz=8, c=GRAY, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "📍 Slide 25 · 项 目 菜 单\n👩‍🏫 老师 选: 一 个 (45 分) 或 两 个 (90 分 — 各 40 分)\n  根 据 时 间 / 材 料\n  K-2 推 荐 Project 1 (转 盘)\n  3-5 都 可 以\n⏱️ 1-2 分 钟 介 绍")


# ============================================================
# 26 · PROJECT 1 STEPS — 垃圾分类转盘
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎡 Project 1 步 骤  ·  Sorting Spinner Steps", MOSS)

tb(s, 0.4, 0.85, 9.2, 0.28, "4 步 做 完 — 一 个 可 以 转 的 分 类 转 盘!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 numbered step cards
p1_steps = [
    ("1️⃣", "🍽️", "拿 一 个 纸 盘",
     "Take a paper plate", MOSS),
    ("2️⃣", "✂️", "分 成 4 个 区",
     "Divide into 4 sections", DEEP_TEAL),
    ("3️⃣", "🖍️", "每 区 画 一 种 垃 圾",
     "Draw a trash type in each", FIRE_ORANGE),
    ("4️⃣", "📌", "中 间 装 一 个 指 针 ",
     "Add a brass-fastener pointer", RECYCLE_BLUE),
]
sw = 2.20; sgap = 0.15
stotal = 4*sw + 3*sgap; sstart = (10 - stotal)/2
for i, (num, em, cn, en, cl) in enumerate(p1_steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 1.30, sw, 3.10, cl, fill=WHITE, lw=2.5)
    # Step number badge
    nb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+sw/2-0.30), Inches(1.40), Inches(0.60), Inches(0.60))
    nb.fill.solid(); nb.fill.fore_color.rgb = cl; nb.line.fill.background()
    tb(s, x+sw/2-0.30, 1.40, 0.60, 0.55, num, sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    # Emoji
    tb(s, x, 2.15, sw, 0.80, em, sz=48, a=PP_ALIGN.CENTER)
    # CN
    tb(s, x+0.05, 3.05, sw-0.10, 0.55, cn, sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
    # EN
    tb(s, x+0.05, 3.62, sw-0.10, 0.42, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    # Arrow between
    if i < 3:
        arrow_x = x + sw + 0.01
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(2.55), Inches(0.13), Inches(0.30))
        arr.fill.solid(); arr.fill.fore_color.rgb = MOSS; arr.line.fill.background()

# Bottom: 4 zones to draw
zb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
zb.fill.solid(); zb.fill.fore_color.rgb = MOSS
zb.line.color.rgb = STAR; zb.line.width = Pt(2.5)
tb(s, 0.55, 4.62, 9.0, 0.30, "🎯 4 个 区 — 画 这 些:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.92, 9.0, 0.42, "📜 纸  ·  🥤 塑 料  ·  🥫 金 属  ·  🍌 食 物 垃 圾",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 26 · Project 1 步 骤\n👩‍🏫 老师 做: 黑板 / 屏幕 示 范 一 遍\n  Step 1-4 都 自 己 做 给 学 生 看\n👧 学生 做: 跟 着 做 (每 步 1 分钟)\n⏱️ 老师 示 范 4-5 分钟")


# ============================================================
# 27 · PROJECT 1 TEACHER DEMO + TIMER
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "⏱️ 动 手 时 间!  ·  Work Time", MOSS)

tb(s, 0.4, 0.85, 9.2, 0.30, "现 在 你 来 做! 老 师 走 动 帮 忙.",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# LEFT: big timer
panel(s, 0.40, 1.20, 4.60, 3.30, MOSS, fill=WHITE, lw=3)
tb(s, 0.40, 1.40, 4.60, 1.50, "⏰", sz=110, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.95, 4.60, 0.55, "30 分 钟",
   sz=42, b=True, c=MOSS, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.55, 4.60, 0.35, "30 minutes · Build time",
   sz=13, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.95, 4.60, 0.32, "🎵 老师 放 一 首 30 分 钟 的 音 乐",
   sz=10, b=True, c=MOSS, a=PP_ALIGN.CENTER)

# RIGHT: teacher walk-around tips
panel(s, 5.20, 1.20, 4.40, 3.30, FIRE_ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.20, 1.20, 4.40, FIRE_ORANGE, "🚶 老 师 走 动 提 醒", sz=13)
walk_tips = [
    "✦ 用 中 文 说: 「这 是 ___ 垃 圾」",
    "✦ 鼓 励 K-2 — 帮 他 们 画 + 写",
    "✦ 3-5 — 多 写 字, 多 说 中 文",
    "✦ 帮 忙 装 指 针 (brass fastener)",
    "✦ 提 醒: 4 个 区 都 要 画 满!",
]
for i, t in enumerate(walk_tips):
    tb(s, 5.35, 1.85 + i*0.50, 4.10, 0.40, t, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)

# Bottom: ready check
rc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.75))
rc.fill.solid(); rc.fill.fore_color.rgb = STAR; rc.line.fill.background()
tb(s, 0.55, 4.75, 9.0, 0.32, "✅ 完 成 检 查: 4 个 区 都 有 图 · 指 针 转 得 动 · 用 中 文 说 出 名 字",
   sz=12, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.22, "Done check: 4 zones drawn · pointer spins · can name in Chinese",
   sz=9, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 27 · 动 手 时 间\n👩‍🏫 老师 做: 启 动 计 时 器 (放 30 分 钟 音 乐 / 用 Google timer)\n  走 动 帮 忙\n👧 学生 做: 安 静 / 小 声 创 作\n⏱️ 30 分 钟\n💡 K-2 需 要 老师 帮 装 brass fastener — 提 前 准 备")


# ============================================================
# 28 · PROJECT 2 — 环保钥匙牌 Shrinky Dink
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔑 Project 2  ·  环 保 钥 匙 牌 · Shrinky Dink Keychain", RECYCLE_BLUE)

tb(s, 0.4, 0.85, 9.2, 0.28, "用 Shrinky Dink 做 一 个 「环 保 小 达 人」 钥 匙 牌!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 4 steps with shrink visualization
p2_steps = [
    ("1️⃣", "🖍️", "在 Shrinky Dink 上 画",
     "Draw on Shrinky Dink"),
    ("2️⃣", "✂️", "剪 出 形 状",
     "Cut out shape"),
    ("3️⃣", "🔥", "老 师 烤 (3-5 分 钟)",
     "Teacher bakes (3-5 min)"),
    ("4️⃣", "🔑", "穿 钥 匙 环 — 完 成!",
     "Add keyring — done!"),
]
sw = 2.20; sgap = 0.15
stotal = 4*sw + 3*sgap; sstart = (10 - stotal)/2
for i, (num, em, cn, en) in enumerate(p2_steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 1.20, sw, 2.20, RECYCLE_BLUE, fill=WHITE, lw=2.5)
    tb(s, x+0.10, 1.30, 0.55, 0.40, num, sz=18, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.LEFT)
    tb(s, x, 1.42, sw, 0.85, em, sz=42, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.35, sw-0.10, 0.45, cn, sz=11, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.82, sw-0.10, 0.40, en, sz=8, c=GRAY, a=PP_ALIGN.CENTER)
    if i < 3:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+sw+0.01), Inches(2.10), Inches(0.13), Inches(0.30))
        arr.fill.solid(); arr.fill.fore_color.rgb = RECYCLE_BLUE; arr.line.fill.background()

# Design ideas
di = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(3.55), Inches(9.20), Inches(1.05))
di.fill.solid(); di.fill.fore_color.rgb = WARM
di.line.color.rgb = RECYCLE_BLUE; di.line.width = Pt(2)
tb(s, 0.55, 3.62, 9.00, 0.32, "🎨 设 计 灵 感  Design Ideas:",
   sz=12, b=True, c=RECYCLE_BLUE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 3.95, 9.00, 0.32, "♻️ Recycle Hero · 🌱 环 保 小 达 人 · 📜 纸 · 🥤 塑 料 · 🥫 金 属 · 🍌 食 物 垃 圾",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.28, 9.00, 0.25, "Eco-Hero badges · 'Recycle Hero' · 4 trash icons",
   sz=9, c=GRAY, a=PP_ALIGN.LEFT)

# Bottom: safety
sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.70))
sb.fill.solid(); sb.fill.fore_color.rgb = FIRE_ORANGE; sb.line.fill.background()
tb(s, 0.55, 4.83, 9.0, 0.30, "⚠️ 安 全: 烤 箱 只 有 老 师 用!  Safety: Only teachers use the oven!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.15, 9.0, 0.22, "学 生 把 画 好 的 交 给 老 师, 老 师 统 一 烘 烤",
   sz=10, b=True, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "📍 Slide 28 · Project 2\n👩‍🏫 老师 做: 提 前 准 备 Shrinky Dink + 油 性 笔 + 钥 匙 环 + 烤 箱\n  Step 1-2 学 生 做, Step 3 老 师 统 一 烤\n👧 学生 做: 画 + 剪 + 等\n⏱️ 学生 25 分 + 老师 烤 5 分\n⚠️ 烤 箱 = 老 师 专 用. 学 生 一 律 不 碰!")


# ============================================================
# 29 · GALLERY WALK + SHARING
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🖼️ 展 示 + 分 享  ·  Gallery Walk & Share", DAY)

tb(s, 0.4, 0.85, 9.2, 0.30, "把 你 的 作 品 摆 在 桌 上 — 大 家 一 起 看 + 说!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)

# 2 sentence frames
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(1.30), Inches(8.80), Inches(1.30))
fr1.fill.solid(); fr1.fill.fore_color.rgb = MOSS
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.42, 8.50, 0.32, "🎡  转 盘 句 型:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.77, 8.50, 0.45, "「这 是 我 的 垃 圾 分 类 转 盘. 我 会 回 收 ___.」",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.28, 8.50, 0.25, "This is my sorting spinner. I will recycle ___.",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)

fr2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(2.80), Inches(8.80), Inches(1.30))
fr2.fill.solid(); fr2.fill.fore_color.rgb = RECYCLE_BLUE
fr2.line.color.rgb = STAR; fr2.line.width = Pt(3)
tb(s, 0.75, 2.92, 8.50, 0.32, "🔑  钥 匙 牌 句 型:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 3.27, 8.50, 0.45, "「这 是 我 的 环 保 钥 匙 牌. 我 要 保 护 地 球!」",
   sz=18, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 3.78, 8.50, 0.25, "This is my eco-keychain. I will protect the Earth!",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)

# Gallery walk how
gw = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.10))
gw.fill.solid(); gw.fill.fore_color.rgb = WARM
gw.line.color.rgb = DAY; gw.line.width = Pt(2.5)
tb(s, 0.55, 4.40, 9.00, 0.32, "🖼️ Gallery Walk  ·  画 廊 散 步:",
   sz=13, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.72, 9.00, 0.30, "1. 作 品 摆 桌 上   2. 全 班 静 静 走 一 圈   3. 喜 欢 的 — 拍 拍 桌 子!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.05, 9.00, 0.30, "Walk silently · tap desk if you love it · share favorites at end",
   sz=10, c=GRAY, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "📍 Slide 29 · Gallery Walk + 分 享\n👩‍🏫 老师 说: 「把 作 品 放 桌 上 — 我 们 静 静 走 一 圈」\n  然 后 抽 4-5 位 学 生 用 句 型 介 绍\n👧 学生 做: 走 + 看 + 说\n⏱️ 10 分钟 (走 5 + 分 享 5)")


# ============================================================
# 30 · SHARE + CLOSE
# ============================================================
s = share_close(prs, DAY,
    frames_cn=["「今 天 我 学 会 了 ______ 是 ______ 垃 圾.」",
               "「我 要 保 护 地 球 — 我 会 ______.」"],
    frames_en="I learned ___ is ___ trash · I will protect the Earth by ___",
    next_day_cn="✨ 明 天: 我 们 来 学 「重 复 使 用 + 减 少 浪 费」!",
    next_day_en="Tomorrow: Reuse & Reduce!",
    next_emoji="🌍")
n += 1; pn(s, n)
notes(s, "📍 Slide 30 · 收 尾\n👩‍🏫 老师 说: 「用 句 型 说 一 句 — 我 学 到 什 么?」\n  抽 4-5 位 学 生\n  全 班 喊: 「我 是 环 保 小 达 人!」\n⏱️ 5 分钟\n🌍 一 天 圆 满 结 束!")


# ============================================================
# Save
# ============================================================
out = os.path.join(os.path.dirname(__file__), "day1_trash.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
