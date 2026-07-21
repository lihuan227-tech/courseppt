#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零废弃 Day 3 — 家庭 Zero Waste 计划  (Family Zero-Waste Plan)
Concept (Zero Waste) + 5R principle (deep on 减少 Reduce + 重复使用 Reuse)
→ hunt for 浪费 at home → design a family eco-action「我们家可以 ____」.

Mascots 熊猫奇奇 & 妙妙 come home with the kids ("来我们家做客").
Adopts the wilderness Day 2 (create_day2_camp.py) helper conventions & style.
Build:  pip install python-pptx  &&  python day3_family_zerowaste.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
W, H = prs.slide_width, prs.slide_height

# --- Palette (zero-waste unit: Forest & Moss) ---
EARTH  = RGBColor(0x2C,0x5F,0x2D)   # primary green (like PINE)
MOSS   = RGBColor(0x6B,0xA0,0x3D)   # secondary green
CREAM  = RGBColor(0xFD,0xF6,0xE3)
WARM   = RGBColor(0xFF,0xF3,0xE0)
SUN    = RGBColor(0xE0,0x7A,0x2C)   # orange accent
SUNYEL = RGBColor(0xF5,0xC2,0x42)
ALERT  = RGBColor(0xC8,0x25,0x3E)
OK     = RGBColor(0x38,0x8E,0x3C)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x2C,0x2C,0x2C)
GRAY   = RGBColor(0x88,0x88,0x8C)
LGRAY  = RGBColor(0xBB,0xBB,0xBB)
IMGBG  = RGBColor(0xE8,0xE8,0xE8)
BROWN  = RGBColor(0x6B,0x44,0x23)
# 5R accent colors (fixed, used consistently)
R_REFUSE = RGBColor(0xC8,0x25,0x3E)  # 拒绝  red
R_REDUCE = RGBColor(0xE0,0x7A,0x2C)  # 减少  orange
R_REUSE  = RGBColor(0x3E,0x6E,0xB6)  # 重复使用 blue
R_RECYC  = RGBColor(0x6B,0xA0,0x3D)  # 回收  green
R_ROT    = RGBColor(0x6B,0x44,0x23)  # 堆肥  brown

FONT = 'KaiTi'

HERE = os.path.dirname(__file__)
ASSET_DIR = os.path.join(HERE, "day3_family_zerowaste_images")
REAL_ASSET_DIR = os.path.join(HERE, "day3_family_zerowaste_real_images")

def _font(size, bold=False):
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for path in candidates:
        if path and os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                pass
    return ImageFont.load_default()

def _rr(draw, xy, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)

def _label(draw, xy, text, size=28, fill=(44, 44, 44), bold=False, anchor="mm"):
    draw.text(xy, text, font=_font(size, bold), fill=fill, anchor=anchor)

def _new_canvas(w, h, bg=(253, 246, 227)):
    return Image.new("RGB", (w, h), bg)

def _save_scene(name, w, h, painter):
    os.makedirs(ASSET_DIR, exist_ok=True)
    path = os.path.join(ASSET_DIR, name)
    img = _new_canvas(w, h)
    draw = ImageDraw.Draw(img)
    painter(draw, w, h)
    img.save(path, quality=95)
    return path

def _draw_panda(draw, x, y, scale=1.0, holding=None):
    s = scale
    draw.ellipse((x-44*s, y-64*s, x+44*s, y+42*s), fill=(255,255,255), outline=(44,44,44), width=max(1,int(3*s)))
    draw.ellipse((x-61*s, y-50*s, x-30*s, y-16*s), fill=(40,40,40))
    draw.ellipse((x+30*s, y-50*s, x+61*s, y-16*s), fill=(40,40,40))
    draw.ellipse((x-24*s, y-30*s, x-6*s, y-8*s), fill=(40,40,40))
    draw.ellipse((x+6*s, y-30*s, x+24*s, y-8*s), fill=(40,40,40))
    draw.ellipse((x-5*s, y-8*s, x+5*s, y+2*s), fill=(40,40,40))
    draw.arc((x-16*s, y-2*s, x, y+16*s), 10, 80, fill=(40,40,40), width=max(1,int(2*s)))
    draw.arc((x, y-2*s, x+16*s, y+16*s), 100, 170, fill=(40,40,40), width=max(1,int(2*s)))
    if holding:
        _label(draw, (x, y+56*s), holding, size=int(32*s), fill=(44,95,45), bold=True)

def _make_assets():
    assets = {}

    def trash_hook(d, w, h):
        d.rectangle((0, int(h*.62), w, h), fill=(221, 236, 206))
        d.rectangle((0, 0, w, int(h*.62)), fill=(207, 232, 236))
        _label(d, (w*.12, h*.12), "家庭垃圾桶", 34, (44,95,45), True, "la")
        d.rectangle((w*.10, h*.36, w*.33, h*.78), fill=(90, 116, 117), outline=(48,72,72), width=5)
        d.rectangle((w*.08, h*.30, w*.35, h*.37), fill=(64, 87, 88), outline=(48,72,72), width=4)
        for bx, by, col in [(0.43,0.47,(230,230,230)),(0.52,0.38,(255,242,179)),(0.60,0.50,(178,218,255)),(0.70,0.42,(235,190,150)),(0.78,0.52,(210,230,170))]:
            _rr(d, (w*bx, h*by, w*(bx+.12), h*(by+.13)), 12, col, (90,90,90), 3)
        d.line((w*.45,h*.76,w*.88,h*.76), fill=(123,92,62), width=6)
        _draw_panda(d, w*.82, h*.35, 0.9, "少一点!")

    def reduce(d, w, h):
        d.rectangle((0, int(h*.70), w, h), fill=(221,236,206))
        _label(d, (w*.08, h*.10), "减少 Reduce", 42, (224,122,44), True, "la")
        for i, (x, y, txt) in enumerate([(0.18,.55,"水壶"),(.44,.56,"两面纸"),(.70,.55,"布袋")]):
            _rr(d, (w*x-70, h*y-80, w*x+70, h*y+90), 28, (255,255,255), (224,122,44), 5)
            if i == 0:
                _rr(d, (w*x-28,h*y-45,w*x+28,h*y+56), 16, (74,151,196), (44,95,45), 4)
                d.rectangle((w*x-18,h*y-67,w*x+18,h*y-45), fill=(44,95,45))
            elif i == 1:
                d.rectangle((w*x-42,h*y-58,w*x+42,h*y+62), fill=(255,252,240), outline=(87,122,70), width=4)
                d.line((w*x-30,h*y-20,w*x+30,h*y-20), fill=(87,122,70), width=3)
                d.line((w*x-30,h*y+8,w*x+30,h*y+8), fill=(87,122,70), width=3)
                d.line((w*x-30,h*y+36,w*x+30,h*y+36), fill=(87,122,70), width=3)
            else:
                d.arc((w*x-34,h*y-78,w*x+34,h*y+5), 180, 360, fill=(44,95,45), width=7)
                _rr(d, (w*x-54,h*y-30,w*x+54,h*y+62), 14, (245,194,66), (44,95,45), 4)
            _label(d, (w*x, h*y+118), txt, 25, (44,44,44), True)
        d.polygon([(w*.48,h*.19),(w*.52,h*.19),(w*.52,h*.31),(w*.56,h*.31),(w*.50,h*.43),(w*.44,h*.31),(w*.48,h*.31)], fill=(224,122,44))

    def reuse(d, w, h):
        d.rectangle((0, int(h*.70), w, h), fill=(221,236,206))
        _label(d, (w*.08, h*.10), "重复使用 Reuse", 40, (62,110,182), True, "la")
        d.arc((w*.18,h*.24,w*.82,h*.88), 20, 320, fill=(62,110,182), width=16)
        d.polygon([(w*.78,h*.31),(w*.91,h*.33),(w*.82,h*.43)], fill=(62,110,182))
        _rr(d, (w*.18,h*.43,w*.34,h*.75), 18, (245,194,66), (44,95,45), 5)
        d.arc((w*.20,h*.35,w*.32,h*.56), 180, 360, fill=(44,95,45), width=5)
        _rr(d, (w*.43,h*.34,w*.55,h*.75), 22, (175,218,231), (62,110,182), 5)
        d.rectangle((w*.46,h*.28,w*.52,h*.34), fill=(62,110,182))
        d.rectangle((w*.66,h*.40,w*.82,h*.73), fill=(207,165,104), outline=(107,68,35), width=5)
        d.line((w*.68,h*.51,w*.80,h*.51), fill=(107,68,35), width=4)
        d.line((w*.68,h*.60,w*.80,h*.60), fill=(107,68,35), width=4)

    def kitchen(d, w, h):
        d.rectangle((0, 0, w, h), fill=(250,238,215))
        d.rectangle((0, h*.70, w, h), fill=(188,211,171))
        d.rectangle((w*.08,h*.17,w*.38,h*.70), fill=(242,244,237), outline=(107,68,35), width=5)
        d.ellipse((w*.17,h*.38,w*.29,h*.48), fill=(200,225,235), outline=(62,110,182), width=3)
        d.line((w*.26,h*.35,w*.34,h*.35), fill=(62,110,182), width=6)
        d.line((w*.32,h*.37,w*.32,h*.60), fill=(62,110,182), width=6)
        for yy in [0.42,0.48,0.54]:
            d.ellipse((w*.30,h*yy,w*.34,h*(yy+.04)), fill=(90,180,220))
        d.rectangle((w*.45,h*.24,w*.82,h*.68), fill=(255,255,255), outline=(107,68,35), width=5)
        d.rectangle((w*.50,h*.44,w*.77,h*.48), fill=(107,68,35))
        d.ellipse((w*.58,h*.33,w*.68,h*.44), fill=(255,245,210), outline=(107,68,35), width=3)
        d.arc((w*.69,h*.27,w*.86,h*.44), 100, 220, fill=(224,122,44), width=6)
        _rr(d, (w*.67,h*.58,w*.88,h*.78), 12, (245,194,66), (44,95,45), 4)
        _label(d, (w*.16,h*.11), "水一直流", 24, (200,37,62), True, "la")
        _label(d, (w*.52,h*.15), "饭菜倒掉", 24, (200,37,62), True, "la")

    def living(d, w, h):
        d.rectangle((0, 0, w, h), fill=(232,240,226))
        d.rectangle((0, h*.72, w, h), fill=(186,150,110))
        d.ellipse((w*.44,h*.05,w*.56,h*.19), fill=(245,194,66), outline=(224,122,44), width=4)
        for a in range(0,360,45):
            pass
        d.line((w*.50,h*.19,w*.50,h*.33), fill=(107,68,35), width=5)
        d.rectangle((w*.12,h*.32,w*.48,h*.62), fill=(60,70,76), outline=(44,44,44), width=5)
        d.rectangle((w*.16,h*.36,w*.44,h*.56), fill=(100,166,205))
        d.rectangle((w*.22,h*.64,w*.38,h*.70), fill=(44,44,44))
        _rr(d, (w*.57,h*.48,w*.86,h*.72), 28, (224,122,44), (107,68,35), 5)
        d.ellipse((w*.64,h*.36,w*.72,h*.52), fill=(175,218,231), outline=(62,110,182), width=4)
        d.rectangle((w*.67,h*.28,w*.70,h*.36), fill=(62,110,182))
        _label(d, (w*.08,h*.13), "没人也开灯", 24, (200,37,62), True, "la")
        _label(d, (w*.10,h*.24), "电视开着", 24, (200,37,62), True, "la")
        _label(d, (w*.58,h*.28), "一次性水瓶", 22, (200,37,62), True, "la")

    def bathroom(d, w, h):
        d.rectangle((0, 0, w, h), fill=(222,242,245))
        d.rectangle((0, h*.70, w, h), fill=(185,214,220))
        d.rectangle((w*.10,h*.18,w*.42,h*.64), fill=(250,250,250), outline=(62,110,182), width=5)
        d.ellipse((w*.18,h*.40,w*.34,h*.53), fill=(205,232,240), outline=(62,110,182), width=4)
        d.line((w*.29,h*.35,w*.38,h*.35), fill=(62,110,182), width=6)
        d.line((w*.36,h*.37,w*.36,h*.58), fill=(62,110,182), width=6)
        for yy in [0.42,0.48,0.54]:
            d.ellipse((w*.34,h*yy,w*.38,h*(yy+.04)), fill=(90,180,220))
        for i in range(4):
            d.ellipse((w*(.53+i*.055),h*.44,w*(.62+i*.055),h*.53), fill=(255,255,255), outline=(180,180,180), width=3)
        d.rectangle((w*.58,h*.54,w*.78,h*.66), fill=(255,255,255), outline=(180,180,180), width=3)
        _rr(d, (w*.78,h*.28,w*.88,h*.65), 18, (160,205,198), (44,137,123), 4)
        d.ellipse((w*.81,h*.20,w*.85,h*.28), fill=(44,137,123))
        _label(d, (w*.11,h*.10), "刷牙水开着", 23, (200,37,62), True, "la")
        _label(d, (w*.52,h*.17), "纸巾太多", 23, (200,37,62), True, "la")

    def poster(d, w, h):
        d.rectangle((0, 0, w, h), fill=(235,244,230))
        d.rectangle((w*.08,h*.08,w*.92,h*.88), fill=(255,252,240), outline=(44,95,45), width=8)
        _label(d, (w*.50,h*.17), "我们家的 Zero Waste 计划", 30, (44,95,45), True)
        rows = [("自带购物袋", (245,194,66)), ("关灯关水", (90,180,220)), ("纸两面用", (107,160,61))]
        for i,(txt,col) in enumerate(rows):
            y = h*(.32+i*.18)
            d.ellipse((w*.16,y-24,w*.24,y+24), fill=col, outline=(44,95,45), width=4)
            d.line((w*.30,y,w*.80,y), fill=(224,122,44), width=5)
            _label(d, (w*.33,y-7), txt, 24, (44,44,44), True, "la")
        _draw_panda(d, w*.77, h*.73, .65, "可以!")

    def class_photo(d, w, h):
        d.rectangle((0, 0, w, h), fill=(207,232,236))
        d.rectangle((0, h*.62, w, h), fill=(221,236,206))
        _label(d, (w*.50,h*.12), "Zero Waste 小卫士", 44, (44,95,45), True)
        xs = [w*.18,w*.32,w*.46,w*.60,w*.74]
        for i,x in enumerate(xs):
            d.ellipse((x-34,h*.32,x+34,h*.50), fill=(255,225,185), outline=(100,70,50), width=3)
            d.arc((x-36,h*.27,x+36,h*.45), 180, 360, fill=(60,45,35), width=14)
            _rr(d, (x-42,h*.52,x+42,h*.80), 18, [(224,122,44),(62,110,182),(107,160,61),(200,37,62),(44,137,123)][i], None, 1)
            d.rectangle((x-65,h*.20,x+65,h*.35), fill=(255,252,240), outline=(44,95,45), width=4)
            _label(d, (x,h*.275), "我们家可以", 17, (44,95,45), True)
        _draw_panda(d, w*.90, h*.56, .75, "合影!")

    assets["trash_hook"] = _save_scene("trash_hook.png", 1290, 780, trash_hook)
    assets["reduce"] = _save_scene("reduce.png", 1290, 990, reduce)
    assets["reuse"] = _save_scene("reuse.png", 1290, 990, reuse)
    assets["kitchen"] = _save_scene("kitchen_waste.png", 1290, 960, kitchen)
    assets["living"] = _save_scene("living_room_waste.png", 1290, 960, living)
    assets["bathroom"] = _save_scene("bathroom_waste.png", 1290, 960, bathroom)
    assets["poster"] = _save_scene("teacher_poster_sample.png", 1380, 1050, poster)
    assets["class_photo"] = _save_scene("class_photo.png", 2100, 750, class_photo)
    return assets

ASSETS = _make_assets()
REAL_ASSETS = {
    "trash_hook": "trash_hook_real.png",
    "reduce": "reduce_real.png",
    "reuse": "reuse_real.png",
    "kitchen": "kitchen_waste_real.png",
    "living": "living_room_waste_real.png",
    "bathroom": "bathroom_waste_real.png",
    "poster": "teacher_poster_sample_real.png",
    "class_photo": "class_photo_real.png",
}
for key, filename in REAL_ASSETS.items():
    real_path = os.path.join(REAL_ASSET_DIR, filename)
    if os.path.exists(real_path):
        ASSETS[key] = real_path
GENERATED_ASSET_DIR = os.path.join(HERE, "assets", "day3_family_zerowaste")
GENERATED_ASSETS = {
    "trash": os.path.join(GENERATED_ASSET_DIR, "trash.png"),
    "reduce": os.path.join(GENERATED_ASSET_DIR, "reduce.png"),
    "reuse": os.path.join(GENERATED_ASSET_DIR, "reuse.png"),
    "kitchen": os.path.join(GENERATED_ASSET_DIR, "kitchen.png"),
    "living": os.path.join(GENERATED_ASSET_DIR, "living_room.png"),
    "bathroom": os.path.join(GENERATED_ASSET_DIR, "bathroom.png"),
    "poster_sample": os.path.join(GENERATED_ASSET_DIR, "poster_sample.png"),
    "class_photo": os.path.join(GENERATED_ASSET_DIR, "class_photo.png"),
}
ASSETS.update({key: path for key, path in GENERATED_ASSETS.items() if os.path.exists(path)})

# === Helpers (same conventions as Day 2 camp) ===
def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=FONT;return tf
def ap(tf,txt,sz=18,b=False,c=DARK,a=None):
    p=tf.add_paragraph()
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=FONT
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def image_card(s,l,t,w,h,path,border_color=MOSS):
    from PIL import Image
    if not os.path.exists(path):
        return False
    with Image.open(path) as im:
        img_ratio=im.width/im.height
    box_ratio=w/h
    pic=s.shapes.add_picture(path,Inches(l),Inches(t),width=Inches(w),height=Inches(h))
    if img_ratio > box_ratio:
        crop=(1-box_ratio/img_ratio)/2
        pic.crop_left=crop; pic.crop_right=crop
    elif img_ratio < box_ratio:
        crop=(1-img_ratio/box_ratio)/2
        pic.crop_top=crop; pic.crop_bottom=crop
    border=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    border.fill.background();border.line.color.rgb=border_color;border.line.width=Pt(2)
    return True
def asset_for_label(lb):
    if "满满一袋垃圾" in lb or "家庭垃圾桶" in lb: return ASSETS["trash"]
    if "减少" in lb: return ASSETS["reduce"]
    if "重复使用" in lb: return ASSETS["reuse"]
    if "海报样板" in lb: return ASSETS["poster_sample"]
    if "全班举着" in lb or "合影" in lb: return ASSETS["class_photo"]
    return None
def ib(s,l,t,w,h,lb="📷"):
    path=asset_for_label(lb)
    if path and image_card(s,l,t,w,h,path,MOSS):
        return
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,lb,sz=14,c=LGRAY,a=PP_ALIGN.CENTER)
def img_box(s,l,t,w,h,path,color=EARTH):
    if image_card(s,l,t,w,h,path,color):
        return
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=IMGBG;sh.line.color.rgb=color;sh.line.width=Pt(2)
def hb(s,txt,c=EARTH,t=0.15):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.55));sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.4,t+0.03,9.2,0.5,txt,sz=20,b=True,c=WHITE)
def pn(s,n): tb(s,9.0,5.25,0.8,0.3,str(n),sz=10,c=GRAY,a=PP_ALIGN.RIGHT)
def notes(s,txt): s.notes_slide.notes_text_frame.text=txt
def div(title,sub,color,emoji=""):
    s=ns();bg(s,color)
    tb(s,1,1.5,8,1.2,f"{emoji} {title}",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,1,2.8,8,0.8,sub,sz=22,c=WHITE,a=PP_ALIGN.CENTER);return s
def pill(s,l,t,w,h,txt,c,sz=14):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,l+0.1,t+h/2-0.2,w-0.2,0.4,txt,sz=sz,b=True,c=WHITE,a=PP_ALIGN.CENTER)
def sentence_frame_bar(s,t,frame_cn,frame_en,label="💬 我来说"):
    sf=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(t),Inches(9.4),Inches(0.65))
    sf.fill.solid();sf.fill.fore_color.rgb=WARM;sf.line.color.rgb=SUN;sf.line.width=Pt(2)
    tb(s,0.5,t+0.1,1.9,0.4,label,sz=14,b=True,c=SUN)
    tb(s,2.3,t+0.07,7.3,0.3,frame_cn,sz=14,b=True,c=DARK)
    tb(s,2.3,t+0.32,7.3,0.3,frame_en,sz=10,c=GRAY)
def plan_frame_bar(s,t):
    """The recurring family-plan frame.「我们家可以 ____。」"""
    sentence_frame_bar(s,t,"我们家可以 ____ 。","Our family can ____ .",label="🏠 家庭计划")

# === Specialized Day 3 helpers ===
def r5_card(s,l,t,w,h,num,cn,en,emoji,color,big=False):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l+0.1),Inches(t+0.08),Inches(0.5),Inches(0.5))
    badge.fill.solid();badge.fill.fore_color.rgb=color;badge.line.fill.background()
    tb(s,l+0.1,t+0.14,0.5,0.4,str(num),sz=16,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,l+0.65,t+0.13,w-0.7,0.35,en,sz=9,c=GRAY)
    tb(s,l+0.05,t+0.70,w-0.1,0.7,emoji,sz=40 if big else 34,a=PP_ALIGN.CENTER)
    tb(s,l+0.05,t+h-0.55,w-0.1,0.45,cn,sz=18 if big else 15,b=True,c=color,a=PP_ALIGN.CENTER)

def r_intro_slide(num,cn,en,emoji,color,idea_cn,idea_en,action_hint,examples_hint):
    """Concept-via-action+picture slide for one R."""
    s=ns();bg(s,CREAM);hb(s,f"{emoji} 5R 第{num}个 · {cn}  {en}",color)
    img_box(s,0.3,1.05,4.3,3.3,ASSETS["reduce" if cn=="减少" else "reuse"],color)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(3.3))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=color;panel.line.width=Pt(2.5)
    head=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(1.05),Inches(4.85),Inches(0.55))
    head.fill.solid();head.fill.fore_color.rgb=color;head.line.fill.background()
    tb(s,5.0,1.13,4.6,0.4,f"{cn} = 什么意思?",sz=15,b=True,c=WHITE)
    tb(s,5.05,1.78,4.5,0.5,idea_cn,sz=18,b=True,c=color)
    tb(s,5.05,2.30,4.5,0.35,idea_en,sz=11,c=GRAY)
    tb(s,5.05,2.78,4.5,0.4,f"🙌 动作: {action_hint}",sz=13,b=True,c=DARK)
    tb(s,5.05,3.25,4.5,1.0,f"💡 例子: {examples_hint}",sz=13,c=DARK)
    plan_frame_bar(s,4.50)
    return s

def example_grid_slide(title,color,examples,frame_kind="plan"):
    """4 example cards in a row + frame bar. examples = [(emoji, cn, en), ...]"""
    s=ns();bg(s,CREAM);hb(s,title,color)
    n_ex=len(examples)
    cw=2.2; gap=(9.4-cw*n_ex)/(n_ex-1) if n_ex>1 else 0
    for i,(em,cn,en) in enumerate(examples):
        x=0.3+i*(cw+gap)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.05),Inches(cw),Inches(2.95))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(2)
        tb(s,x+0.05,1.30,cw-0.1,0.9,em,sz=46,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,2.45,cw-0.1,0.5,cn,sz=15,b=True,c=color,a=PP_ALIGN.CENTER)
        tb(s,x+0.05,3.05,cw-0.1,0.7,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    plan_frame_bar(s,4.30)
    return s

def r_light_slide(num,cn,en,emoji,color,meaning,examples_line,callback=None):
    """Light-touch R: one slide, big card + examples + (optional) callback note."""
    s=ns();bg(s,CREAM);hb(s,f"{emoji} 5R 第{num}个 · {cn}  {en}",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.15),Inches(4.3),Inches(3.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    tb(s,0.5,1.45,4.1,1.1,emoji,sz=84,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.95,4.1,0.6,cn,sz=30,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.6,4.1,0.4,en,sz=14,c=GRAY,a=PP_ALIGN.CENTER)
    panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.95),Inches(1.15),Inches(4.75),Inches(3.0))
    panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=color;panel.line.width=Pt(2)
    tf=tb(s,5.15,1.35,4.4,0.6,f"意思: {meaning}",sz=16,b=True,c=color)
    ap(tf,"",sz=6)
    ap(tf,f"💡 例子: {examples_line}",sz=14,c=DARK)
    if callback:
        ap(tf,"",sz=6)
        ap(tf,f"🔁 {callback}",sz=12,c=GRAY)
    plan_frame_bar(s,4.40)
    return s

def match_slide(title,subtitle,items):
    """配一配: behavior → which R. Direct-answer (the R is shown on the card)."""
    s=ns();bg(s,CREAM);hb(s,title,EARTH)
    tb(s,0.4,0.80,9.2,0.30,subtitle,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
    n_it=len(items)
    cw=2.95; gap=(9.4-cw*n_it)/(n_it-1) if n_it>1 else 0
    for i,(em,beh_cn,beh_en,r_cn,r_color) in enumerate(items):
        x=0.3+i*(cw+gap)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.25),Inches(cw),Inches(2.65))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=r_color;sh.line.width=Pt(2)
        tb(s,x+0.05,1.45,cw-0.1,0.8,em,sz=44,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,2.45,cw-0.2,0.5,beh_cn,sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,3.0,cw-0.2,0.35,beh_en,sz=9,c=GRAY,a=PP_ALIGN.CENTER)
        pill(s,x+cw/2-0.95,3.40,1.9,0.40,f"→ {r_cn}",r_color,sz=13)
    plan_frame_bar(s,4.30)
    return s

def find_waste_slide(room_cn,room_en,color,wastes):
    """家庭找浪费: scene image LEFT + waste list (问题→怎么办) RIGHT + plan frame."""
    s=ns();bg(s,CREAM);hb(s,f"🔍 家里找浪费 · {room_cn}  {room_en}",color)
    tb(s,0.4,0.80,9.2,0.28,"圈出浪费 → 想一想:「我们家可以 ____ 。」 Circle the waste, then say how we fix it.",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    room_key={"厨房":"kitchen","客厅":"living","浴室":"bathroom"}[room_cn]
    img_box(s,0.4,1.20,4.30,3.20,ASSETS[room_key],color)
    y0=1.20
    h_each=(3.20-(len(wastes)-1)*0.12)/len(wastes)
    for i,(em,what_cn,fix_cn) in enumerate(wastes):
        y=y0+i*(h_each+0.12)
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.85),Inches(y),Inches(4.85),Inches(h_each))
        sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=color;sh.line.width=Pt(1.5)
        tb(s,4.95,y+(h_each-0.5)/2,0.55,0.5,em,sz=24,a=PP_ALIGN.CENTER)
        tb(s,5.55,y+0.08,4.05,0.35,f"❌ {what_cn}",sz=13,b=True,c=ALERT)
        tb(s,5.55,y+0.42,4.05,0.35,f"✅ 我们可以 {fix_cn}",sz=13,b=True,c=OK)
    plan_frame_bar(s,4.50)
    return s

def word_card_read(w,py,en,sent,img,color=SUN):
    s=ns();bg(s,CREAM);hb(s,"👀 我会认  I Can Read",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.5),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
    wsz={1:72,2:60,3:46,4:36}.get(len(w),34)
    tb(s,0.5,1.15,4.3,1.3,w,sz=wsz,b=True,c=EARTH,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.45,4.3,0.4,f"{py}  {en}",sz=18,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.5,2.88,4.3,0.4,"👉 跟我读！Read after me!",sz=14,c=color,a=PP_ALIGN.CENTER)
    ib(s,5.3,1.0,4.4,2.5,img)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.8),Inches(9.2),Inches(1.2))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,0.6,3.9,1.5,0.4,"例句 Example",sz=14,b=True,c=color)
    tb(s,0.6,4.3,8.8,0.5,sent,sz=20,b=True,c=DARK)
    return s

def word_card_write(w,py,en,strokes_hint,color=EARTH):
    # day4_emergency format: big-char card + 3 步练习 + 田字格 + teacher/student bar
    s=ns();bg(s,CREAM);hb(s,f"✍️ 我会写 · {w}  I Can Write",color)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.0),Inches(4.4),Inches(3.4))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=color;sh.line.width=Pt(2.5)
    csz={1:150,2:96,3:72,4:56}.get(len(w),56)
    tb(s,0.5,1.2,4.2,1.9,w,sz=csz,b=True,c=color,a=PP_ALIGN.CENTER)
    tb(s,0.5,3.1,4.2,0.45,f"{py}  ·  {en}",sz=18,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,0.55,3.62,4.1,0.45,strokes_hint,sz=10,c=DARK,a=PP_ALIGN.CENTER)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.0),Inches(4.6),Inches(1.6))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=color;sh2.line.width=Pt(2)
    tb(s,5.15,1.1,4.4,0.4,"✏️ 3 步练习  3 Steps",sz=16,b=True,c=color)
    tb(s,5.15,1.55,4.4,0.35,"1️⃣ 看老师写  Watch teacher",sz=13,c=DARK)
    tb(s,5.15,1.90,4.4,0.35,"2️⃣ 用手指空中写  Air-write",sz=13,c=DARK)
    tb(s,5.15,2.25,4.4,0.35,"3️⃣ 在田字格写 3 次",sz=13,c=DARK)
    for i in range(4):
        x=5.0+i*1.15
        sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.85),Inches(1.05),Inches(1.05))
        sq.fill.solid();sq.fill.fore_color.rgb=WHITE;sq.line.color.rgb=color;sq.line.width=Pt(1.5)
        ln1=s.shapes.add_connector(1,Inches(x),Inches(3.375),Inches(x+1.05),Inches(3.375));ln1.line.color.rgb=LGRAY;ln1.line.width=Pt(0.5);ln1.line.dash_style=2
        ln2=s.shapes.add_connector(1,Inches(x+0.525),Inches(2.85),Inches(x+0.525),Inches(3.9));ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.5);ln2.line.dash_style=2
    tb(s,5.0,3.95,4.6,0.3,"在田字格里写 3 次 ↓",sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.5),Inches(9.2),Inches(0.62))
    bar.fill.solid();bar.fill.fore_color.rgb=WARM;bar.line.color.rgb=color;bar.line.width=Pt(1.5)
    tb(s,0.55,4.55,4.4,0.24,"👩‍🏫 老师问 Teacher asks:",sz=10,b=True,c=color)
    tb(s,0.55,4.80,4.4,0.26,f"和我一起写「{w}」",sz=12,b=True,c=DARK)
    tb(s,5.05,4.55,4.5,0.24,"🧒 学生 Student does:",sz=10,b=True,c=color)
    tb(s,5.05,4.80,4.5,0.26,"看 → 空中写 → 田字格写 3 次",sz=12,b=True,c=DARK)
    return s

# ========================================================================
#                              SLIDES
# ========================================================================
n=0

# 1. COVER
s=ns();bg(s,EARTH)
sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.4),W,Inches(2.0))
sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.fill.background()
tb(s,1,0.4,8,0.5,"DAY 3",sz=18,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1,0.95,8,0.7,"🐼🏠 家庭 Zero Waste 计划",sz=38,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.7,8,0.5,"Our Family's Zero-Waste Plan",sz=20,c=WARM,a=PP_ALIGN.CENTER)
tb(s,1,2.6,8,0.5,"🌱 奇奇妙妙来我们家做客！",sz=24,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,1,3.15,8,0.4,"5R · 找浪费 · 我们家可以 ____",sz=14,b=True,c=MOSS,a=PP_ALIGN.CENTER)
tb(s,1,3.55,8,0.4,"Refuse · Reduce · Reuse · Recycle · Rot",sz=13,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,1,4.6,8,0.4,"零废弃 · Zero Waste",sz=14,b=True,c=SUN,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"开场 (1 分钟):\n• 「小朋友们好！还记得熊猫奇奇和妙妙吗？今天它们要来我们家做客！」\n• 它们想帮我们家做一个「Zero Waste 零废弃」计划。\n• 今天 3 件大事: 1) 什么是零废弃 + 5R; 2) 在家里找浪费; 3) 做我们家的计划 + 环保购物袋。\n• 关键句型 (一整天都用): 「我们家可以 ____ 。」")

# 2. SESSION 1 DIVIDER
s=div("Session 1  上午","🌍 什么是 Zero Waste? + 5R 大家庭",EARTH,"📖");n+=1;pn(s,n)

# 3. LEARNING GOALS
s=ns();bg(s,CREAM);hb(s,"🎯 今天的目标  Today's Goals",EARTH)
goals=[("🌍","懂什么是「零废弃」","Understand Zero Waste"),
       ("♻️","学会 5R 环保原则","Learn the 5R principle"),
       ("🔍","在家里找出浪费","Spot waste at home"),
       ("🏠","做我们家的环保计划","Make our family plan")]
y=1.15
for em,cn,en in goals:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(y),Inches(8.8),Inches(0.85))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MOSS;sh.line.width=Pt(2)
    tb(s,0.75,y+0.12,0.8,0.6,em,sz=30,a=PP_ALIGN.CENTER)
    tb(s,1.7,y+0.12,6.0,0.5,cn,sz=20,b=True,c=EARTH)
    tb(s,1.7,y+0.52,6.0,0.3,en,sz=11,c=GRAY)
    y+=0.98
n+=1;pn(s,n)
notes(s,"读目标 (1 分钟): 一条一条读, 让学生跟读。重点是第 4 条 — 今天每个人都要做一个「我们家的计划」带回家。")

# 4. HOOK — mascots visit + how much trash
s=ns();bg(s,CREAM);hb(s,"🐼 奇奇妙妙来做客  Our Pandas Visit",EARTH)
tb(s,0.4,0.85,9.2,0.45,"🗑️ 想一想: 我们家一天会扔多少垃圾?",sz=24,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,0.4,1.35,9.2,0.30,"How much trash does our family throw away in one day?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
img_box(s,0.5,1.85,4.3,2.6,ASSETS["trash_hook"],MOSS)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.0),Inches(1.85),Inches(4.6),Inches(2.6))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=MOSS;panel.line.width=Pt(2.5)
tf=tb(s,5.2,2.05,4.3,0.5,"🐼 奇奇说:",sz=15,b=True,c=EARTH)
ap(tf,"我们扔的垃圾, 大部分不会消失,",sz=14,c=DARK)
ap(tf,"会埋在地下很久很久。",sz=14,c=DARK)
ap(tf,"",sz=6)
ap(tf,"🐼 妙妙说:",sz=15,b=True,c=SUN)
ap(tf,"如果少扔一点, 地球就会更干净!",sz=14,c=DARK)
sentence_frame_bar(s,4.60,"我们家一天大约扔 ___ 袋垃圾。","Our family throws about ___ bag(s) a day.")
n+=1;pn(s,n)
notes(s,"激趣 (3 分钟):\n• 举起一袋垃圾 (或图片): 「这是一个家一天的垃圾, 多吗?」\n• 让学生猜自己家一天扔几袋。\n• 关键概念: 垃圾不会消失, 会埋在地下很久。少产生垃圾 = 帮助地球。\n• 引出今天的大问题: 「我们能不能少扔一点? 怎么做?」")

# 5. WHAT IS ZERO WASTE
s=ns();bg(s,CREAM);hb(s,"🌍 什么是 Zero Waste?  零废弃",EARTH)
tb(s,0.4,0.85,9.2,0.5,"零废弃 = 尽量少产生垃圾, 让东西不浪费。",sz=22,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,0.4,1.40,9.2,0.30,"Zero Waste = make as little trash as possible; don't waste things.",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
cards=[("🚫🗑️","少产生垃圾","Make less trash",MOSS),
       ("🔄","东西重复用","Use things again",R_REUSE),
       ("🌱","爱护地球","Care for Earth",EARTH)]
for i,(em,cn,en,cl) in enumerate(cards):
    x=0.6+i*3.0
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.95),Inches(2.7),Inches(2.0))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.15,2.6,0.8,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.15,2.6,0.5,cn,sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.6,0.3,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.25,"零废弃就是 ___ 。","Zero waste means ___.")
n+=1;pn(s,n)
notes(s,"讲概念 (3 分钟) — 用动作 + 图片:\n• 「零废弃」三个字拆开: 零 = 没有, 废弃 = 扔掉的东西。合起来 = 尽量不扔东西。\n• 三个动作带孩子做: 🚫(摇手)少产生垃圾 / 🔄(转手)东西重复用 / 🌱(抱)爱护地球。\n• 不是「一点垃圾都没有」, 而是「尽量少」。让孩子明白这是一个努力的方向。")

# 6. 浪费是什么 (recognize 浪费)
s=ns();bg(s,CREAM);hb(s,"🍽️ 浪费是什么?  What is 浪费 (Waste)?",SUN)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.95),Inches(4.5),Inches(2.3))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.fill.background()
tb(s,0.5,1.05,4.3,1.0,"浪费",sz=60,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,0.5,2.25,4.3,0.4,"làng fèi  ·  waste",sz=18,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.5,2.75,4.3,0.4,"= 还能用却扔掉 / 用太多",sz=15,b=True,c=DARK,a=PP_ALIGN.CENTER)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(0.95),Inches(4.6),Inches(2.3))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=SUN;panel.line.width=Pt(2)
tf=tb(s,5.3,1.10,4.3,0.45,"哪些是浪费? Which are waste?",sz=14,b=True,c=SUN)
ap(tf,"💧 水开着一直流",sz=14,c=DARK)
ap(tf,"🍚 饭吃不完倒掉",sz=14,c=DARK)
ap(tf,"💡 没人也开着灯",sz=14,c=DARK)
ap(tf,"📄 纸只用一面就扔",sz=14,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(3.4),Inches(9.3),Inches(0.7))
sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=OK;sh2.line.width=Pt(2)
tb(s,0.55,3.52,9.0,0.45,"👀 我会认: 浪费  —  跟我读三遍! 浪 · 费 · 浪费!",sz=16,b=True,c=OK)
sentence_frame_bar(s,4.25,"这是浪费, 因为 ___ 。","This is waste, because ___.")
n+=1;pn(s,n)
notes(s,"教生词「浪费」(3 分钟):\n• 出示词卡「浪费」, 跟读三遍。\n• 用具体例子让孩子判断「这是不是浪费?」: 水一直流 / 倒饭 / 开着灯 / 只用一面的纸。\n• 让孩子用句型说: 「这是浪费, 因为 还能用 / 太多了。」\n• 「浪费」是今天的核心词 — 找浪费、不浪费。")

# 7. THINK-PAIR-SHARE
s=ns();bg(s,CREAM);hb(s,"🤔 想一想 · 说一说  Think–Pair–Share",MOSS)
tb(s,0.4,0.95,9.2,0.5,"在我们家, 你看到过什么浪费?",sz=24,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,0.4,1.5,9.2,0.3,"What waste have you seen at home?",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
steps=[("🧠","想","Think","自己先想一个"),
       ("🗣️","说","Share","和同桌说一说"),
       ("👂","听","Listen","听别人怎么说")]
for i,(em,cn,en,desc) in enumerate(steps):
    x=0.8+i*2.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.5),Inches(1.9))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MOSS;sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.15,2.4,0.7,em,sz=40,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.95,2.4,0.45,f"{cn} {en}",sz=17,b=True,c=EARTH,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.45,2.4,0.35,desc,sz=12,c=DARK,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.15,"我看到过 ___ 浪费。","I have seen ___ wasted.")
n+=1;pn(s,n)
notes(s,"Think-Pair-Share (3-4 分钟):\n• 想 (30 秒): 自己想一个家里看到的浪费。\n• 说: 和同桌互相说。老师计时 1 分钟。\n• 听: 请 2-3 个孩子分享, 用句型「我看到过 ___ 浪费。」\n• 老师把答案记在白板上, 留到 Session 2 找浪费时呼应。")

# 8. 5R LINEUP
s=ns();bg(s,CREAM);hb(s,"♻️ 5R 大家庭  The 5R Family",EARTH)
tb(s,0.4,0.80,9.2,0.32,"5 个环保好习惯 — 从上往下, 越前面越棒!  Best habits first ↓",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
r5=[(1,"拒绝","Refuse","🙅",R_REFUSE),
    (2,"减少","Reduce","⬇️",R_REDUCE),
    (3,"重复使用","Reuse","🔄",R_REUSE),
    (4,"回收","Recycle","♻️",R_RECYC),
    (5,"堆肥","Rot","🍂",R_ROT)]
cw=1.78; gap=(9.4-cw*5)/4
for i,(num,cn,en,em,cl) in enumerate(r5):
    x=0.3+i*(cw+gap)
    big=cn in("减少","重复使用")
    r5_card(s,x,1.25,cw,2.55,num,cn,en,em,cl,big=big)
tb(s,0.4,3.95,9.2,0.35,"⭐ 今天重点学: 减少 (Reduce) + 重复使用 (Reuse)",sz=15,b=True,c=SUN,a=PP_ALIGN.CENTER)
plan_frame_bar(s,4.45)
n+=1;pn(s,n)
notes(s,"介绍 5R (4 分钟):\n• 一个一个介绍, 让学生跟读 + 做动作: 拒绝🙅(摆手) / 减少⬇️(往下压) / 重复使用🔄(转圈) / 回收♻️(投桶) / 堆肥🍂(撒).\n• 顺序有意义: 越前面越好 — 最好的是「拒绝」(一开始就不要), 最后才是「回收/堆肥」。\n• 今天重点是中间两个: 减少 + 重复使用 (黄色加大的卡)。\n• 拒绝/回收/堆肥 简单认识即可。回收可以呼应 Day 1 的四个垃圾桶。")

# 9. 减少 INTRO
s=r_intro_slide(2,"减少","Reduce","⬇️",R_REDUCE,
    "用得更少, 不要太多。","Use less; don't take too much.",
    "双手往下压 = 减少 ⬇️","少用一次性的东西、买需要的就好")
n+=1;pn(s,n)
notes(s,"教「减少 Reduce」(3 分钟):\n• 动作: 双手像往下压, 说「减——少」。\n• 意思: 一开始就用得少。例: 自带水壶不买瓶装水、纸两面用、买大包装少买小包装。\n• 「减少」是今天要写的字之一, 多读几遍。\n• 句型: 「我们家可以 减少 ___ 。」")

# 10. 减少 EXAMPLES
s=example_grid_slide("⬇️ 减少 Reduce · 我们家可以这样做",R_REDUCE,
    [("🥤","自带水壶","Bring a bottle"),
     ("📄","纸两面用","Use both sides"),
     ("🛍️","自带购物袋","Bring own bag"),
     ("🍚","吃多少盛多少","Take what you eat")])
n+=1;pn(s,n)
notes(s,"减少的例子 (2-3 分钟):\n• 一张一张看, 让孩子说「这样可以减少什么?」(减少瓶子 / 减少用纸 / 减少塑料袋 / 减少剩饭)。\n• 用句型: 「我们家可以 自带水壶 / 纸两面用 / ……」\n• 问: 「你们家做过哪一个?」举手。")

# 11. 重复使用 INTRO
s=r_intro_slide(3,"重复使用","Reuse","🔄",R_REUSE,
    "用过的东西再用一次。","Use something again instead of throwing it.",
    "双手转圈圈 = 重复使用 🔄","旧瓶子、旧衣服、购物袋都能再用")
n+=1;pn(s,n)
notes(s,"教「重复使用 Reuse」(3 分钟):\n• 动作: 双手转圈, 说「重复——使用」。\n• 意思: 东西没坏, 再用一次, 不要马上扔。\n• 例: 玻璃瓶装东西、旧衣服给弟弟妹妹、纸箱做手工、购物袋下次再带。\n• 重点带出「购物袋」这个生词 (下一页例子里有)。")

# 12. 重复使用 EXAMPLES (购物袋)
s=example_grid_slide("🔄 重复使用 Reuse · 我们家可以这样做",R_REUSE,
    [("🛍️","购物袋再用","Reuse the bag"),
     ("🫙","旧瓶子装东西","Reuse jars"),
     ("👕","旧衣服传给弟妹","Hand down clothes"),
     ("📦","纸箱做手工","Craft from boxes")])
n+=1;pn(s,n)
notes(s,"重复使用的例子 (2-3 分钟):\n• 第一张「购物袋」— 教生词「购物袋」: 去超市买东西装的袋子。自己带, 就不用每次拿新的塑料袋。\n• 其它例子让孩子说一说自己家有没有做。\n• 句型: 「我们家可以 重复使用 ___ 。」")

# 13. 拒绝 (light)
s=r_light_slide(1,"拒绝","Refuse","🙅",R_REFUSE,
    "不需要的, 一开始就说「不用了」。",
    "不要免费小玩具 · 不要塑料吸管 · 不要多余的袋子")
n+=1;pn(s,n)
notes(s,"简单认识「拒绝 Refuse」(1-2 分钟):\n• 动作: 摆摆手说「不用了, 谢谢!」\n• 最厉害的一招 — 一开始就不拿, 就不会变成垃圾。\n• 例: 餐厅的塑料吸管、免费的小传单、用不到的赠品。")

# 14. 回收 (light, callback Day 1)
s=r_light_slide(4,"回收","Recycle","♻️",R_RECYC,
    "可回收的东西, 送去做成新东西。",
    "塑料瓶 · 纸 · 玻璃 · 金属 → 蓝色可回收桶",
    callback="还记得 Day 1 的四个桶吗? 可回收物放蓝桶!")
n+=1;pn(s,n)
notes(s,"简单认识「回收 Recycle」(1-2 分钟):\n• 呼应 Day 1: 我们学过四个垃圾桶, 可回收物放蓝桶。\n• 回收 = 让旧东西变成新东西 (旧瓶子 → 新瓶子)。\n• 但回收要花很多力气, 所以「减少」「重复使用」更好, 排在回收前面。")

# 15. 堆肥 (light)
s=r_light_slide(5,"堆肥","Rot","🍂",R_ROT,
    "果皮、菜叶埋在土里, 变成肥料。",
    "果皮 · 菜叶 · 落叶 → 堆肥 → 给植物当营养")
n+=1;pn(s,n)
notes(s,"简单认识「堆肥 Rot」(1-2 分钟):\n• 厨房的果皮菜叶不一定是垃圾 — 埋进土里会烂掉, 变成肥料, 给植物吃。\n• 呼应 Day 1 的厨余垃圾 (绿桶)。\n• 「让垃圾变成有用的东西!」")

# 16. MATCH 1
s=match_slide("🧩 配一配 (1)  Which R?","看行为, 说一说它是哪个 R — 直接看答案!",
    [("🛍️","自带购物袋","Bring own bag","重复使用",R_REUSE),
     ("🥤","自带水壶","Bring a bottle","减少",R_REDUCE),
     ("🥤🚫","不要塑料吸管","No plastic straw","拒绝",R_REFUSE)])
n+=1;pn(s,n)
notes(s,"配一配 (2-3 分钟) — 直接给答案式:\n• 出示行为, 全班一起说是哪个 R, 然后看卡片上的答案对一对。\n• 自带购物袋 = 重复使用; 自带水壶 = 减少; 不要吸管 = 拒绝。\n• 有的行为可能同时是「减少 + 重复使用」, 都对, 鼓励孩子说理由。")

# 17. MATCH 2
s=match_slide("🧩 配一配 (2)  Which R?","看行为, 说一说它是哪个 R — 直接看答案!",
    [("🍂","果皮变肥料","Peels → compost","堆肥",R_ROT),
     ("📄","纸两面用","Use both sides","减少",R_REDUCE),
     ("♻️","旧报纸放蓝桶","Paper → blue bin","回收",R_RECYC)])
n+=1;pn(s,n)
notes(s,"配一配 (2-3 分钟):\n• 果皮变肥料 = 堆肥; 纸两面用 = 减少; 旧报纸放蓝桶 = 回收。\n• 复习全部 5R。问: 「哪个最好? 为什么?」(拒绝/减少最好 — 一开始就不产生垃圾)。")

# 18. S1 EXIT TICKET
s=ns();bg(s,CREAM);hb(s,"🎟️ 出门票  Exit Ticket",EARTH)
tb(s,0.4,1.0,9.2,0.6,"说一句话再下课!",sz=26,b=True,c=EARTH,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.2),Inches(1.9),Inches(7.6),Inches(1.5))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
tb(s,1.4,2.2,7.2,0.6,"「我们家可以 ____ 。」",sz=30,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1.4,2.95,7.2,0.35,"Our family can ____ .",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,3.7,9.2,0.4,"⬇️ 减少 · 🔄 重复使用 · 🙅 拒绝 · ♻️ 回收 · 🍂 堆肥",sz=16,b=True,c=EARTH,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"出门票 (2 分钟): 每个孩子用句型说一件事才能下课 (或排队出门)。\n• 「我们家可以 自带购物袋 / 关灯 / 纸两面用……」\n• 鼓励用上一个 5R 的词。下午继续在家里找浪费!")

# ======================= SESSION 2 =======================
# 19. S2 DIVIDER
s=div("Session 2  下午","🔍 家庭找浪费 + 词汇 + 书写",MOSS,"🔍");n+=1;pn(s,n)

# 20-22. FIND WASTE x3
s=find_waste_slide("厨房","Kitchen",R_REDUCE,
    [("💧","水龙头一直开着","随手关水龙头"),
     ("🍚","饭菜吃不完倒掉","吃多少盛多少"),
     ("🛍️","一次性塑料袋","自带购物袋")])
n+=1;pn(s,n)
notes(s,"找浪费 · 厨房 (3 分钟):\n• 看图, 让孩子先「圈出」浪费的地方 (口头指 / 投影上圈)。\n• 一个一个对: 水一直流、倒饭、用一次性袋子。\n• 每个浪费配一个办法, 用句型「我们家可以 关水 / 吃多少盛多少 / 自带购物袋」。")

s=find_waste_slide("客厅","Living Room",R_REUSE,
    [("💡","没人还开着灯","离开就关灯"),
     ("📺","没人看的电视","不看就关电视"),
     ("🥤","一次性水瓶","用自己的水壶")])
n+=1;pn(s,n)
notes(s,"找浪费 · 客厅 (3 分钟):\n• 灯没人开着、电视没人看、用一次性水瓶。\n• 句型「我们家可以 关灯 / 关电视 / 用水壶」。\n• 关灯关电视也是省电 — 呼应 Day 2 能源。")

s=find_waste_slide("浴室","Bathroom",R_RECYC,
    [("🚿","刷牙时水一直流","刷牙时关水"),
     ("🧻","用太多纸巾","用手帕 / 少抽几张"),
     ("🧴","沐浴露挤太多","用多少挤多少")])
n+=1;pn(s,n)
notes(s,"找浪费 · 浴室 (3 分钟):\n• 刷牙时水一直流、用太多纸巾、沐浴露挤太多。\n• 句型「我们家可以 刷牙关水 / 用手帕 / 用多少挤多少」。\n• 总结三个房间: 家里到处都能找到浪费, 只要留心就能改!")

# 23. 我会认 review (拍词卡)
s=ns();bg(s,CREAM);hb(s,"👀 我会认 · 拍词卡  Slap the Word",SUN)
tb(s,0.4,0.85,9.2,0.35,"老师说词, 谁先拍到谁赢! Teacher says it — slap it fast!",sz=14,b=True,c=DARK,a=PP_ALIGN.CENTER)
words=[("减少","jiǎn shǎo"),("重复使用","chóng fù shǐ yòng"),("环保","huán bǎo"),
       ("购物袋","gòu wù dài"),("浪费","làng fèi")]
cw=1.78; gap=(9.4-cw*5)/4
for i,(w,py) in enumerate(words):
    x=0.3+i*(cw+gap)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.5),Inches(cw),Inches(2.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
    wsz={2:30,3:26,4:20}.get(len(w),20)
    tb(s,x+0.03,1.95,cw-0.06,0.9,w,sz=wsz,b=True,c=EARTH,a=PP_ALIGN.CENTER)
    tb(s,x+0.03,2.95,cw-0.06,0.6,py,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,3.95,9.2,0.35,"💡 我会认 5 个词: 减少 · 重复使用 · 环保 · 购物袋 · 浪费",sz=14,b=True,c=EARTH,a=PP_ALIGN.CENTER)
plan_frame_bar(s,4.45)
n+=1;pn(s,n)
notes(s,"拍词卡 (3-4 分钟):\n• 把 5 个词卡贴白板, 老师说一个 (中文或英文), 学生上来拍。\n• 或全班分两组比赛。\n• 每个词读 2-3 遍, 配动作: 减少⬇️ / 重复使用🔄 / 环保🌱 / 购物袋🛍️ / 浪费🍽️。")

# 24. 词汇配对
s=ns();bg(s,CREAM);hb(s,"🔗 词汇配对  Match the Word",MOSS)
tb(s,0.4,0.82,9.2,0.32,"把词和图连起来! Match each word to its picture.",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
left=[("减少","A"),("重复使用","B"),("环保","C"),("购物袋","D"),("浪费","E")]
right=[("⬇️ 用得更少","2"),("🛍️ 装东西的袋子","5"),("🍽️ 还能用却扔掉","4"),("🌱 爱护地球","3"),("🔄 再用一次","1")]
y=1.25
for (w,lab),(pic,num) in zip(left,right):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(4.2),Inches(0.62))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MOSS;sh.line.width=Pt(2)
    pill(s,0.6,y+0.11,0.45,0.4,lab,MOSS,sz=13)
    tb(s,1.2,y+0.13,3.4,0.4,w,sz=18,b=True,c=EARTH)
    sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.3),Inches(y),Inches(4.2),Inches(0.62))
    sh2.fill.solid();sh2.fill.fore_color.rgb=WHITE;sh2.line.color.rgb=SUN;sh2.line.width=Pt(2)
    pill(s,5.4,y+0.11,0.45,0.4,num,SUN,sz=13)
    tb(s,6.0,y+0.15,3.4,0.4,pic,sz=15,b=True,c=DARK)
    y+=0.74
n+=1;pn(s,n)
notes(s,"词汇配对 (3 分钟):\n• 左边词, 右边意思, 打乱了。让孩子把字母和数字连起来。\n• 答案: 减少=2(用得更少), 重复使用=1(再用一次), 环保=3(爱护地球), 购物袋=5(装东西的袋子), 浪费=4(还能用却扔掉)。\n• 可以做成连线练习纸。")

# 25. 句型练习
s=ns();bg(s,CREAM);hb(s,"🗣️ 句型练习  Sentence Practice",EARTH)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.05),Inches(8.0),Inches(1.2))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(3)
tb(s,1.2,1.35,7.6,0.6,"「我们家可以 ____ 。」",sz=32,b=True,c=SUN,a=PP_ALIGN.CENTER)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(2.5),Inches(8.0),Inches(2.2))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=MOSS;panel.line.width=Pt(2)
tf=tb(s,1.25,2.65,7.5,0.45,"🐼 跟同桌一问一答: Pair talk —",sz=14,b=True,c=EARTH)
ap(tf,"A: 我们家可以做什么环保的事?",sz=17,b=True,c=R_REUSE)
ap(tf,"B: 我们家可以 自带购物袋 。",sz=17,b=True,c=SUN)
ap(tf,"A: 我们家可以 关灯 。我们家可以 纸两面用 。",sz=15,c=DARK)
n+=1;pn(s,n)
notes(s,"句型练习 (4 分钟):\n• 全班齐读句型 3 遍, 拍手打节奏。\n• 同桌一问一答, 每人至少说 2 句「我们家可以 ____」。\n• 鼓励用上 5R 的词和今天的生词 (减少/重复使用/购物袋)。\n• 这个句型下午做海报时还要用。")

# 26-27. 我会写
s=word_card_write("减少","jiǎn shǎo","reduce",
    "「减」: 左边三点水 氵, 右边「咸」。\n「少」: 先写「小」, 再加一撇 丿。\n两个字都不太难, 慢慢写。",color=R_REDUCE)
n+=1;pn(s,n)
notes(s,"写「减少」(4 分钟):\n• 先看老师在田字格写一遍, 边写边说部件。\n• 「减」= 氵 + 咸; 「少」= 小 + 丿。\n• 空中书空 → 田字格写 3 遍。\n• 提醒占格: 字要写在田字格中间, 不要太大或太小。")

s=word_card_write("环保","huán bǎo","eco-friendly",
    "「环」: 左边「王(玉)」旁, 右边「不」。\n「保」: 左边单人旁 亻, 右边「呆」。\n环保 = 爱护环境。",color=EARTH)
n+=1;pn(s,n)
notes(s,"写「环保」(4 分钟):\n• 「环」= 王字旁 + 不; 「保」= 单人旁 + 呆。\n• 老师范写 → 书空 → 田字格写 3 遍。\n• 连起来读「环保」, 并说一句「我要做环保小卫士!」")

# 28. BAMBOOZLE
s=ns();bg(s,MOSS)
tb(s,1,0.55,8,0.7,"🎮 复习游戏  Review Game",sz=34,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.3,8,0.5,"Bamboozle · 零废弃大挑战",sz=20,b=True,c=WARM,a=PP_ALIGN.CENTER)
card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.5),Inches(2.1),Inches(7.0),Inches(1.7))
card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.fill.background()
tf=tb(s,1.8,2.3,6.4,0.5,"👉 点击打开 / Teacher opens:",sz=15,b=True,c=EARTH)
ap(tf,"www.baamboozle.com  (老师课前 import 题目)",sz=14,c=DARK)
ap(tf,"📄 题库文件: bamboozle_day3_family.csv  (同一文件夹, ~10 题)",sz=13,c=GRAY)
tb(s,1,4.05,8,0.4,"题目涵盖: 5R · 找浪费 · 生词 (减少/重复使用/环保/购物袋/浪费)",sz=13,b=True,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"复习游戏 (8-10 分钟):\n• 课前在 baamboozle.com 用 bamboozle_day3_family.csv 建好题目 (~10 题)。\n• 分 2-4 组, 答对加分。答错全班一起说正确答案。\n• 题目涵盖 5R、找浪费、5 个生词。")

# 28b. REVIEW: BAMBOOZLE (separate review page, blank link)
s=ns();bg(s,CREAM);hb(s,"🎮 复习 · Bamboozle 大复习  Review",MOSS)
tb(s,0.4,0.85,9.2,0.3,"复习游戏 — 分组抢答!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
lp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(1.2),Inches(4.4),Inches(3.6))
lp.fill.solid();lp.fill.fore_color.rgb=EARTH;lp.line.fill.background()
tb(s,0.4,1.5,4.4,1.0,"🎮",sz=72,a=PP_ALIGN.CENTER)
btn=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.05),Inches(2.8),Inches(3.1),Inches(0.6))
btn.fill.solid();btn.fill.fore_color.rgb=SUN;btn.line.fill.background()
tb(s,1.05,2.9,3.1,0.42,"▶️ 开始复习",sz=15,b=True,c=WHITE,a=PP_ALIGN.CENTER)
lk=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.7),Inches(3.95),Inches(3.8),Inches(0.55))
lk.fill.solid();lk.fill.fore_color.rgb=WHITE;lk.line.color.rgb=SUNYEL;lk.line.width=Pt(1.5)
tb(s,0.85,4.05,3.55,0.38,"🔗 链接: ____________ (老师粘贴)",sz=11,b=True,c=GRAY,a=PP_ALIGN.LEFT)
rp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.1),Inches(1.2),Inches(4.5),Inches(3.6))
rp.fill.solid();rp.fill.fore_color.rgb=WHITE;rp.line.color.rgb=SUN;rp.line.width=Pt(2.5)
tb(s,5.3,1.35,4.1,0.4,"📋 怎么玩 How to Play",sz=15,b=True,c=SUN)
for i,t in enumerate(["1. 老师点开上面的链接","2. 全班分 2-3 组","3. 轮流抢答, 答对加分","4. 复习今天: 5R + 找浪费 + 生词"]):
    tb(s,5.35,1.95+i*0.6,4.1,0.5,t,sz=13,b=True,c=DARK)
n+=1;pn(s,n)
notes(s,"复习游戏 — 链接留空, 老师课前粘贴 Bamboozle 链接")

# 29. REFLECTION
s=ns();bg(s,CREAM);hb(s,"💭 今天我学会了  I Learned",EARTH)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.0),Inches(1.3),Inches(8.0),Inches(1.4))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=SUN;sh.line.width=Pt(2.5)
tb(s,1.2,1.6,7.6,0.6,"「今天我学会了 ____ 。」",sz=28,b=True,c=SUN,a=PP_ALIGN.CENTER)
tb(s,1.2,2.25,7.6,0.35,"Today I learned ____ .",sz=14,c=GRAY,a=PP_ALIGN.CENTER)
tb(s,0.4,3.15,9.2,0.4,"🌱 明天: 做我们家的 Zero Waste 计划 + 环保购物袋!",sz=18,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,0.4,3.75,9.2,0.35,"Next: make our family plan + decorate a reusable bag!",sz=12,c=GRAY,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"小结 (2 分钟): 请几个孩子说「今天我学会了 ___」。\n• 预告下午/明天: 做海报 + 装饰环保购物袋。\n• 布置: 回家观察一个浪费, 想一个「我们家可以 ___」。")

# ======================= SESSION 3 =======================
# 30. S3 DIVIDER
s=div("Session 3  动手时间","🛠️ 我们家的 Zero Waste 计划 + 环保购物袋",SUN,"🛠️");n+=1;pn(s,n)

# 30b. COMPLETE BOOKLET (before projects)
s=ns();bg(s,CREAM);hb(s,"📓 先完成练习册!  Booklet First",EARTH)
tb(s,0.4,0.85,9.2,0.3,"做手工之前, 先完成练习册!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
books=[("Day 1","垃圾分类",R_RECYC),("Day 2","可再生能源",SUN),("Day 3","水与塑料",R_REUSE),("Day 4","家庭计划",MOSS)]
bw=2.15;bgap=0.27;bstart=(10-4*bw-3*bgap)/2
for i,(lab,zh,cl) in enumerate(books):
    x=bstart+i*(bw+bgap)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.3),Inches(bw),Inches(3.1))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x,1.55,bw,0.85,"📓",sz=46,a=PP_ALIGN.CENTER)
    tb(s,x,2.5,bw,0.45,f"{lab}",sz=18,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x,2.95,bw,0.35,zh,sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
    cb=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+bw/2-0.85),Inches(3.55),Inches(1.7),Inches(0.45))
    cb.fill.solid();cb.fill.fore_color.rgb=WARM;cb.line.color.rgb=cl;cb.line.width=Pt(1.2)
    tb(s,x+bw/2-0.85,3.62,1.7,0.34,"✅ 完成",sz=11,b=True,c=OK,a=PP_ALIGN.CENTER)
tbar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(4.7),Inches(9.2),Inches(0.55))
tbar.fill.solid();tbar.fill.fore_color.rgb=WARM;tbar.line.color.rgb=SUNYEL;tbar.line.width=Pt(1.5)
tb(s,0.55,4.8,9.0,0.36,"👩‍🏫 先看 → 一起读 → 自己写 → 同桌检查",sz=12,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"做手工之前先完成练习册 (Day 1-4)")

# 31. PROJECT MENU
s=ns();bg(s,CREAM);hb(s,"🛠️ 两个动手活动  Two Hands-On Projects",EARTH)
projects=[("活动 1 · 核心","🏠 家庭 Zero Waste 计划海报","Family Plan Poster",
           "从 5R 里选 3 件事, 全家一起做","🏠",MOSS),
          ("活动 2","🛍️ 装饰环保购物袋","Decorate a Reusable Bag",
           "画上图案, 带去超市重复使用","🛍️",R_REUSE)]
for i,(lbl,nm,en,desc,em,cl) in enumerate(projects):
    x=0.4+i*4.6
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(0.95),Inches(4.4),Inches(4.15))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2.5)
    tb(s,x+0.1,1.05,4.2,0.35,lbl,sz=12,b=True,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.4,4.2,0.55,nm,sz=19,b=True,c=cl,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,1.98,4.2,0.35,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.1,2.45,4.2,1.0,em,sz=60,a=PP_ALIGN.CENTER)
    tb(s,x+0.2,3.75,4.0,0.7,desc,sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"介绍两个活动 (2 分钟):\n• 活动 1 (核心, 每人都做): 家庭 Zero Waste 计划海报 — 从 5R 选 3 件事, 写「我们家可以 ___」。\n• 活动 2 (动手): 装饰一个环保购物袋, 带回家重复使用。\n• 90 分钟: 海报 ~30 分钟 + 购物袋 ~25 分钟 + 分享/合影。")

# 32. POSTER — materials + steps
s=ns();bg(s,CREAM);hb(s,"🏠 活动1: 家庭计划海报 · 材料 + 做法",MOSS)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.95),Inches(4.3),Inches(0.4))
sh.fill.solid();sh.fill.fore_color.rgb=MOSS;sh.line.fill.background()
tb(s,0.4,0.98,4.1,0.35,"🧺 材料  Materials",sz=14,b=True,c=WHITE)
tf=tb(s,0.4,1.45,4.3,2.3,"📄 海报纸 / 卡纸",sz=14,c=DARK)
ap(tf,"🖍️ 彩笔 / 蜡笔",sz=14,c=DARK)
ap(tf,"🏷️ 5R 小图标贴纸 (可选)",sz=14,c=DARK)
ap(tf,"✏️「我们家可以 ___」模板纸",sz=14,c=DARK)
sh2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.9),Inches(0.95),Inches(4.8),Inches(0.4))
sh2.fill.solid();sh2.fill.fore_color.rgb=SUN;sh2.line.fill.background()
tb(s,5.0,0.98,4.6,0.35,"👉 做法  Steps",sz=14,b=True,c=WHITE)
tf2=tb(s,5.0,1.45,4.7,2.3,"1️⃣ 写上标题: 我们家的 Zero Waste 计划",sz=13,c=DARK)
ap(tf2,"2️⃣ 从 5R 里选 3 件能做的事",sz=13,c=DARK)
ap(tf2,"3️⃣ 每件写一句「我们家可以 ___」",sz=13,c=DARK)
ap(tf2,"4️⃣ 画图 / 贴图标, 装饰漂亮",sz=13,c=DARK)
sh3=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(3.95),Inches(9.4),Inches(1.15))
sh3.fill.solid();sh3.fill.fore_color.rgb=WARM;sh3.line.color.rgb=MOSS;sh3.line.width=Pt(2)
tb(s,0.5,4.05,9,0.35,"🗣️ 海报上写这样的句子:",sz=14,b=True,c=MOSS)
tb(s,0.5,4.45,4.5,0.35,"· 我们家可以 自带购物袋 。",sz=14,c=DARK)
tb(s,0.5,4.75,4.5,0.35,"· 我们家可以 关灯关水 。",sz=14,c=DARK)
tb(s,5.1,4.45,4.5,0.35,"· 我们家可以 纸两面用 。",sz=14,c=DARK)
tb(s,5.1,4.75,4.5,0.35,"· 我们家可以 剩饭少一点 。",sz=14,c=DARK)
n+=1;pn(s,n)
notes(s,"海报材料 + 做法 (讲 2 分钟):\n• 每人一张海报纸, 写「我们家的 Zero Waste 计划」。\n• 核心要求: 从 5R 选 3 件家里能做的事, 每件写一句「我们家可以 ___」。\n• 字不会写的可以画图 / 用图标贴纸代替。重点是想法, 不是字漂亮。")

# 33. POSTER — 想一想 (choose 3 from 5R)
s=ns();bg(s,CREAM);hb(s,"🤔 想一想: 从 5R 选 3 件事  Choose 3",MOSS)
tb(s,0.4,0.80,9.2,0.30,"哪 3 件事我们家最容易做到? Which 3 can our family really do?",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
ideas=[("⬇️","减少","自带水壶 · 纸两面用 · 少买一次性",R_REDUCE),
       ("🔄","重复使用","购物袋再用 · 旧瓶装东西 · 旧衣传给弟妹",R_REUSE),
       ("🙅","拒绝","不要塑料吸管 · 不要多余的袋子",R_REFUSE),
       ("♻️","回收","瓶子/纸/玻璃放蓝桶",R_RECYC),
       ("🍂","堆肥","果皮菜叶变肥料",R_ROT)]
y=1.20
for em,cn,exs,cl in ideas:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(y),Inches(9.0),Inches(0.62))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.65,y+0.10,0.6,0.45,em,sz=24,a=PP_ALIGN.CENTER)
    tb(s,1.35,y+0.13,2.0,0.4,cn,sz=17,b=True,c=cl)
    tb(s,3.4,y+0.16,6.0,0.35,exs,sz=13,c=DARK)
    y+=0.72
plan_frame_bar(s,4.45)
n+=1;pn(s,n)
notes(s,"想一想 (3 分钟) — 在动手前先选好:\n• 带孩子过一遍 5R 的点子, 让每个孩子心里选 3 件「我们家最容易做到的」。\n• 提醒: 选真的能做到的, 回家要和爸爸妈妈一起做。\n• 选好就开始写海报。")

# 34. POSTER — teacher demo
s=ns();bg(s,CREAM);hb(s,"👩‍🏫 老师示范  Teacher Demo",MOSS)
img_box(s,0.4,1.0,4.6,3.5,ASSETS["poster"],MOSS)
panel=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.2),Inches(1.0),Inches(4.5),Inches(3.5))
panel.fill.solid();panel.fill.fore_color.rgb=WHITE;panel.line.color.rgb=MOSS;panel.line.width=Pt(2.5)
tf=tb(s,5.4,1.2,4.1,0.5,"📋 一张好海报有:",sz=16,b=True,c=MOSS)
ap(tf,"✅ 大标题: 我们家的 Zero Waste 计划",sz=14,c=DARK)
ap(tf,"✅ 3 句「我们家可以 ___」",sz=14,c=DARK)
ap(tf,"✅ 每句配一个小图 / 图标",sz=14,c=DARK)
ap(tf,"✅ 涂上颜色, 漂漂亮亮",sz=14,c=DARK)
ap(tf,"",sz=8)
ap(tf,"⏱️ 等下有 30 分钟做!",sz=14,b=True,c=SUN)
n+=1;pn(s,n)
notes(s,"示范 (2-3 分钟):\n• 老师出示自己做好的样板海报, 边指边说三句「我们家可以 ___」。\n• 强调结构: 标题 + 3 句 + 配图 + 上色。\n• 不求完美, 求想法清楚、能讲出来。")

# 35. POSTER — work time
s=ns();bg(s,MOSS)
tb(s,1,0.7,8,0.8,"✏️ 动手做海报!",sz=40,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.6,8,0.5,"Make Your Family Plan Poster",sz=20,c=WARM,a=PP_ALIGN.CENTER)
clock=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.8),Inches(2.3),Inches(2.4),Inches(2.4))
clock.fill.solid();clock.fill.fore_color.rgb=WHITE;clock.line.color.rgb=SUN;clock.line.width=Pt(4)
tb(s,3.8,2.75,2.4,0.8,"⏱️",sz=54,a=PP_ALIGN.CENTER)
tb(s,3.8,3.65,2.4,0.6,"30 分钟",sz=26,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"画完记得想一想: 怎么和同学讲我的计划?",sz=14,b=True,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"做海报 (30 分钟):\n• 巡视, 帮助选不出 3 件事的孩子 (提示 5R)。\n• 鼓励大孩子多写字, 小孩子多画图。\n• 提前 5 分钟提醒收尾, 准备装饰购物袋。\n• 做完先想好怎么向同学介绍。")

# 36. BAG — materials
s=ns();bg(s,CREAM);hb(s,"🛍️ 活动2: 环保购物袋 · 材料  Materials",R_REUSE)
tb(s,0.4,0.85,9.2,0.35,"为什么做? 自己的袋子重复使用 = 减少塑料袋!",sz=15,b=True,c=R_REUSE,a=PP_ALIGN.CENTER)
mats=[("👜","布袋 / 纸袋","Cloth or paper bag (1/人)"),
      ("🖍️","布用彩笔 / 蜡笔","Fabric markers / crayons"),
      ("🧩","图案模板 (可选)","Stencils (optional)"),
      ("📰","垫在袋子里的纸","Paper to slip inside")]
for i,(em,cn,en) in enumerate(mats):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.45),Inches(2.2),Inches(2.5))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=R_REUSE;sh.line.width=Pt(2)
    tb(s,x+0.05,1.7,2.1,0.8,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.75,2.1,0.5,cn,sz=15,b=True,c=R_REUSE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.35,2.1,0.5,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.30,"这是我的购物袋, 我可以重复使用。","This is my bag — I can reuse it!")
n+=1;pn(s,n)
notes(s,"购物袋材料 (1 分钟):\n• 每人一个布袋或纸袋。布袋用布用彩笔, 纸袋用蜡笔/马克笔。\n• 袋子里垫张纸, 防止颜色透到背面。\n• 老师强调: 这个袋子带回家, 下次买东西重复使用 — 这就是 Reuse!")

# 37. BAG — 4 steps
s=ns();bg(s,CREAM);hb(s,"🛍️ 环保购物袋 · 做法 4 步  Steps",R_REUSE)
steps=[("1️⃣","袋里垫纸","Slip paper inside","防止透色"),
       ("2️⃣","先用铅笔画","Sketch with pencil","轻轻画草图"),
       ("3️⃣","彩笔涂颜色","Color it in","画环保图案"),
       ("4️⃣","晾干 + 写名字","Dry + name it","写上自己名字")]
for i,(num,cn,en,desc) in enumerate(steps):
    x=0.4+i*2.35
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.15),Inches(2.2),Inches(2.7))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=R_REUSE;sh.line.width=Pt(2)
    tb(s,x+0.05,1.3,2.1,0.6,num,sz=34,b=True,c=R_REUSE,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.0,2.1,0.6,cn,sz=17,b=True,c=DARK,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,2.6,2.1,0.4,en,sz=10,c=GRAY,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.1,2.1,0.6,desc,sz=12,c=R_REUSE,a=PP_ALIGN.CENTER)
sentence_frame_bar(s,4.30,"我在袋子上画了 ___ 。","I drew ___ on my bag.")
n+=1;pn(s,n)
notes(s,"购物袋做法 (1 分钟):\n• 4 步: 垫纸 → 铅笔打草稿 → 彩笔上色 → 晾干写名字。\n• 提醒先用铅笔轻轻画, 画错可以改。\n• 下一页给图案点子。")

# 38. BAG — design ideas
s=ns();bg(s,CREAM);hb(s,"🎨 画什么? 图案点子  Design Ideas",R_REUSE)
tb(s,0.4,0.85,9.2,0.32,"画上和环保有关的图案吧! Draw something green!",sz=13,b=True,c=DARK,a=PP_ALIGN.CENTER)
ideas=[("🌳","大树 / 森林"),("🌍","地球"),("♻️","回收标志"),
       ("🌸","花和叶子"),("🐼","奇奇妙妙"),("☀️","太阳 / 蓝天")]
for i,(em,cn) in enumerate(ideas):
    row=i//3; col=i%3
    x=0.7+col*3.05; y=1.35+row*1.55
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.85),Inches(1.35))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MOSS;sh.line.width=Pt(2)
    tb(s,x+0.05,y+0.1,2.75,0.7,em,sz=40,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,y+0.92,2.75,0.35,cn,sz=15,b=True,c=EARTH,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"图案点子 (1 分钟):\n• 给犹豫的孩子一些点子: 大树、地球、回收标志、花叶、熊猫、太阳。\n• 也可以写字: 环保、Reuse、我爱地球。\n• 鼓励原创, 不用都一样。")

# 39. BAG — work time + safety
s=ns();bg(s,R_REUSE)
tb(s,1,0.7,8,0.8,"🎨 动手装饰购物袋!",sz=36,b=True,c=WHITE,a=PP_ALIGN.CENTER)
clock=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.8),Inches(1.75),Inches(2.0),Inches(2.0))
clock.fill.solid();clock.fill.fore_color.rgb=WHITE;clock.line.color.rgb=SUN;clock.line.width=Pt(4)
tb(s,3.8,2.1,2.0,0.7,"⏱️",sz=46,a=PP_ALIGN.CENTER)
tb(s,3.8,2.9,2.0,0.5,"25 分钟",sz=24,b=True,c=R_REUSE,a=PP_ALIGN.CENTER)
safe=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(1.2),Inches(4.05),Inches(7.6),Inches(1.1))
safe.fill.solid();safe.fill.fore_color.rgb=WARM;safe.line.color.rgb=ALERT;safe.line.width=Pt(2)
tf=tb(s,1.4,4.15,7.2,0.4,"🛡️ 安全 Safety:",sz=14,b=True,c=ALERT)
ap(tf,"彩笔只画在袋子上 · 不放进嘴里 · 画完洗手",sz=13,b=True,c=DARK)
n+=1;pn(s,n)
notes(s,"装饰购物袋 (25 分钟):\n• 巡视, 提醒先铅笔后彩笔。\n• 布用彩笔可能有味道, 提醒不要靠近嘴巴和眼睛, 画完洗手。\n• 提前 5 分钟收尾、晾干、写名字。")

# 40. GALLERY WALK
s=ns();bg(s,CREAM);hb(s,"🖼️ 作品展览  Gallery Walk",EARTH)
tb(s,0.4,0.95,9.2,0.5,"摆出作品, 安静走一圈, 看看别人的!",sz=22,b=True,c=EARTH,a=PP_ALIGN.CENTER)
steps=[("🪑","摆桌","Display","把海报和袋子摆好"),
       ("🚶","静走","Walk quietly","安静走一圈看"),
       ("👏","拍桌","Tap & cheer","喜欢就轻拍桌子")]
for i,(em,cn,en,desc) in enumerate(steps):
    x=0.8+i*2.85
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.5),Inches(2.2))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=MOSS;sh.line.width=Pt(2.5)
    tb(s,x+0.05,2.2,2.4,0.8,em,sz=44,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.15,2.4,0.45,f"{cn} {en}",sz=16,b=True,c=EARTH,a=PP_ALIGN.CENTER)
    tb(s,x+0.05,3.65,2.4,0.45,desc,sz=12,c=DARK,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"展览 (3-4 分钟):\n• 学生把海报 + 购物袋摆在桌上。\n• 安静走一圈, 看别人的作品, 喜欢就轻拍桌子 (不出声评价)。\n• 老师挑 2-3 个有创意的, 等下请上台分享。")

# 41. SHARING — poster frame
s=ns();bg(s,CREAM);hb(s,"🗣️ 分享 (1): 我的家庭计划  Share Poster",MOSS)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(1.1),Inches(8.8),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=MOSS;sh.line.width=Pt(2.5)
tf=tb(s,0.85,1.35,8.3,0.5,"举起海报, 这样说:",sz=16,b=True,c=MOSS)
ap(tf,"· 这是我们家的 Zero Waste 计划。",sz=20,b=True,c=DARK)
ap(tf,"· 我们家可以 ___ 。",sz=20,b=True,c=SUN)
ap(tf,"· 我们家可以 ___ 。",sz=20,b=True,c=SUN)
ap(tf,"· 我们家可以 ___ 。",sz=20,b=True,c=SUN)
n+=1;pn(s,n)
notes(s,"分享海报 (4-5 分钟):\n• 请几个孩子上台举起海报, 用句型说出 3 件「我们家可以 ___」。\n• K 级说 1 句即可, G1-3 说 3 句。\n• 全班给掌声。")

# 42. SHARING — bag frame
s=ns();bg(s,CREAM);hb(s,"🗣️ 分享 (2): 我的环保购物袋  Share Bag",R_REUSE)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(1.1),Inches(8.8),Inches(3.0))
sh.fill.solid();sh.fill.fore_color.rgb=WARM;sh.line.color.rgb=R_REUSE;sh.line.width=Pt(2.5)
tf=tb(s,0.85,1.35,8.3,0.5,"举起购物袋, 这样说:",sz=16,b=True,c=R_REUSE)
ap(tf,"· 这是我的环保购物袋。",sz=20,b=True,c=DARK)
ap(tf,"· 我在上面画了 ___ 。",sz=20,b=True,c=SUN)
ap(tf,"· 我可以重复使用它去买东西。",sz=20,b=True,c=EARTH)
n+=1;pn(s,n)
notes(s,"分享购物袋 (3-4 分钟):\n• 请几个孩子举起袋子, 说画了什么 + 「我可以重复使用它」。\n• 强调 Reuse 的意义: 带这个袋子去超市, 就少用一个塑料袋。")

# 43. GROUP PHOTO
s=ns();bg(s,EARTH)
tb(s,1,0.6,8,0.8,"📸 全班合影!",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1,1.5,8,0.5,"Class Photo — hold up your work!",sz=18,c=WARM,a=PP_ALIGN.CENTER)
img_box(s,1.5,2.2,7.0,2.5,ASSETS["class_photo"],WHITE)
tb(s,1,4.95,8,0.4,"🐼 我们都是 Zero Waste 小卫士!",sz=18,b=True,c=SUNYEL,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"合影 (2 分钟): 全班举着海报和购物袋拍照。\n• 喊口号: 「我们家可以 — 零废弃!」\n• 照片可以发给家长, 鼓励回家真的执行计划。")

# ======================= CLOSE =======================
# 44. 家园共育 / 在家练习
s=ns();bg(s,CREAM);hb(s,"🏡 在家练习 · 家园共育  At-Home Practice",EARTH)
tb(s,0.4,0.85,9.2,0.4,"把今天的计划带回家, 全家一起做!",sz=20,b=True,c=EARTH,a=PP_ALIGN.CENTER)
items=[("📌","把海报贴在家里","贴在冰箱 / 门口, 天天看得见",MOSS),
       ("👨‍👩‍👧","全家一起做一件事","这周先做到 1 件「我们家可以 ___」",SUN),
       ("🛍️","带上环保购物袋","下次买东西重复使用",R_REUSE),
       ("🔍","继续找浪费","谁发现浪费, 谁来提醒大家",R_REDUCE)]
y=1.45
for em,cn,desc,cl in items:
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.6),Inches(y),Inches(8.8),Inches(0.8))
    sh.fill.solid();sh.fill.fore_color.rgb=WHITE;sh.line.color.rgb=cl;sh.line.width=Pt(2)
    tb(s,0.75,y+0.18,0.8,0.5,em,sz=28,a=PP_ALIGN.CENTER)
    tb(s,1.7,y+0.1,3.3,0.5,cn,sz=17,b=True,c=cl)
    tb(s,5.0,y+0.22,4.3,0.4,desc,sz=13,c=DARK)
    y+=0.9
n+=1;pn(s,n)
notes(s,"家园共育 (2 分钟):\n• 海报贴在家里显眼处 (冰箱/门口)。\n• 这周先认真做到 1 件事, 不用一下子全做。\n• 购物袋随身带。\n• 鼓励孩子当「监督员」, 发现家人浪费就友好提醒。\n• 可拍照/打卡反馈给老师。")

# 45. SHARE + CLOSE
s=ns();bg(s,EARTH)
tb(s,1,0.7,8,0.8,"🌱 我们家可以…",sz=42,b=True,c=WHITE,a=PP_ALIGN.CENTER)
sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(3.6),Inches(1.9),Inches(2.8),Inches(2.8))
sh.fill.solid();sh.fill.fore_color.rgb=SUNYEL;sh.line.color.rgb=WHITE;sh.line.width=Pt(4)
tb(s,3.6,2.35,2.8,0.7,"🐼",sz=72,a=PP_ALIGN.CENTER)
tb(s,3.6,3.5,2.8,0.5,"零废弃小卫士",sz=20,b=True,c=EARTH,a=PP_ALIGN.CENTER)
tb(s,3.6,4.05,2.8,0.4,"Zero-Waste Hero",sz=12,c=BROWN,a=PP_ALIGN.CENTER)
tb(s,1,5.0,8,0.4,"谢谢奇奇妙妙! 我们家可以一起保护地球 🌍",sz=15,b=True,c=WARM,a=PP_ALIGN.CENTER)
n+=1;pn(s,n)
notes(s,"结束 (1 分钟):\n• 全班一起说一遍「我们家可以 ____」。\n• 谢谢熊猫奇奇和妙妙来做客。\n• 「每个人做一点点, 地球就会很不一样!」\n• 颁发「零废弃小卫士」贴纸 / 印章。")

# === Save ===
out="/Users/Huan/0 projects/summercourse/Chinese/zero_waste零废弃/day3_family_zerowaste.pptx"
prs.save(out)
print(f"Saved {out}  ({n} slides)")
