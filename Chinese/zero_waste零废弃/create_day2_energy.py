#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零 废 弃 与 可 持 续 发 展 · Day 2: 能 量 从 哪 里 来?  (可 再 生 能 源)
小 小 能 源 工 程 师  ·  Little Energy Engineers
Follows day1_trash_v2 / day2_camp style. 3-session deck, K-5 Chinese immersion.

Per spec: docs/superpowers/specs/2026-05-31-zero-waste-day2-energy-design.md
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from _helpers import *
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = make_presentation()
DAY = SKY          # clean-energy blue
SOLAR = STAR       # ☀️ 太阳能
WIND = SKY         # 💨 风能
HYDRO = DEEP_TEAL  # 💧 水力
FOSSIL = INK       # 🛢️ 化石能源 (runs out)
n = 0

HERE = os.path.dirname(__file__)
CAR_DIR = os.path.join(HERE, "互动版_太阳能小车_科学区_大班", "太阳能小车_科学区_大班_步骤图")
GEN_DIR = os.path.join(HERE, "互动版_太阳能发电_科学区_大班", "太阳能发电_科学区_大班_步骤图")
CAR_VIDEO = os.path.join(HERE, "互动版_太阳能小车_科学区_大班", "太阳能小车_科学区_大班_操作视频.mp4")
GEN_VIDEO = os.path.join(HERE, "互动版_太阳能发电_科学区_大班", "太阳能发电_科学区_大班_操作视频.mp4")


def arrow(s, x, y, w=0.20, h=0.30, color=DAY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


def img_or_slot(s, l, t, w, h, path, sug_cn, sug_en, color):
    """Insert a real image if it exists, else a photo placeholder."""
    if path and os.path.exists(path):
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = IMGBG
        box.line.color.rgb = color; box.line.width = Pt(2)
        s.shapes.add_picture(path, Inches(l+0.08), Inches(t+0.08), Inches(w-0.16), Inches(h-0.16))
    else:
        photo_slot(s, l, t, w, h, sug_cn, sug_en, color)


# ============================================================
# DATA — energies + practice + vocab
# ============================================================
RENEWABLES = [
    {
        "emoji": "☀️", "name_cn": "太 阳 能", "name_en": "Solar Energy",
        "color": SOLAR,
        "examples": [
            ("✅", "太 阳 能 板 — 晒 太 阳 发 电", "Solar panels make electricity"),
            ("✅", "太 阳 能 热 水 器", "Solar water heaters"),
            ("✅", "太 阳 能 小 车 / 计 算 器", "Solar toy cars, calculators"),
            ("🌞", "太 阳 每 天 都 来 — 用 不 完!", "The sun comes every day — never runs out!"),
            ("💡", "白 天 有 太 阳 才 行", "Needs daytime sunshine"),
        ],
        "frame_cn": "「太 阳 能 来 自 太 阳.」",
        "frame_en": "Solar energy comes from the sun.",
    },
    {
        "emoji": "💨", "name_cn": "风 能", "name_en": "Wind Energy",
        "color": WIND,
        "examples": [
            ("✅", "风 车 / 风 力 发 电 机", "Wind turbines make electricity"),
            ("✅", "帆 船 — 风 推 着 走", "Sailboats pushed by wind"),
            ("✅", "风 筝 — 风 让 它 飞", "Kites fly on wind"),
            ("🌬️", "风 一 直 吹 — 用 不 完!", "Wind keeps blowing — never runs out!"),
            ("💡", "有 风 的 时 候 才 行", "Needs windy weather"),
        ],
        "frame_cn": "「风 能 来 自 风.」",
        "frame_en": "Wind energy comes from wind.",
    },
    {
        "emoji": "💧", "name_cn": "水 力", "name_en": "Hydro Energy",
        "color": HYDRO,
        "examples": [
            ("✅", "水 坝 / 水 电 站 发 电", "Dams make electricity"),
            ("✅", "水 车 — 水 流 推 动", "Waterwheels turned by flowing water"),
            ("✅", "瀑 布 的 大 力 量", "Powerful waterfalls"),
            ("💧", "水 一 直 流 — 用 不 完!", "Water keeps flowing — never runs out!"),
            ("💡", "要 有 流 动 的 水", "Needs moving water"),
        ],
        "frame_cn": "「水 力 来 自 流 动 的 水.」",
        "frame_en": "Hydro energy comes from flowing water.",
    },
]

# 5 practice items — picture → which energy?  (direct-answer panels)
PRACTICE = [
    {
        "title": "☀️ 太 阳 能 板", "title_en": "Solar panel",
        "img_label": "📸 屋 顶 上 的 太 阳 能 板", "color": SOLAR,
        "panels": [
            {"mark": "✅", "q": "这 是 什 么 能 源?",
             "lines": ["太 阳 能! 它 来 自 太 阳",
                       "Solar — it comes from the sun",
                       "晒 太 阳 就 能 发 电 · Sunlight → electricity"]},
            {"mark": "🌞", "q": "会 用 完 吗?",
             "lines": ["不 会! 太 阳 每 天 都 来",
                       "No — the sun comes every day",
                       "用 不 完 的 能 源 · Renewable!"]},
        ],
        "frame_cn": "「太 阳 能 来 自 太 阳.」",
        "frame_en": "Solar energy comes from the sun.",
    },
    {
        "title": "💨 风 力 发 电 机", "title_en": "Wind turbine",
        "img_label": "📸 大 大 的 白 色 风 车", "color": WIND,
        "panels": [
            {"mark": "✅", "q": "这 是 什 么 能 源?",
             "lines": ["风 能! 它 来 自 风",
                       "Wind — it comes from the wind",
                       "风 吹 → 风 车 转 → 发 电"]},
            {"mark": "🌬️", "q": "会 用 完 吗?",
             "lines": ["不 会! 风 一 直 吹",
                       "No — wind keeps blowing",
                       "用 不 完 的 能 源 · Renewable!"]},
        ],
        "frame_cn": "「风 能 来 自 风.」",
        "frame_en": "Wind energy comes from wind.",
    },
    {
        "title": "💧 水 坝", "title_en": "Hydro dam",
        "img_label": "📸 大 水 坝 + 流 水", "color": HYDRO,
        "panels": [
            {"mark": "✅", "q": "这 是 什 么 能 源?",
             "lines": ["水 力! 它 来 自 流 动 的 水",
                       "Hydro — it comes from flowing water",
                       "水 流 → 推 动 机 器 → 发 电"]},
            {"mark": "💧", "q": "会 用 完 吗?",
             "lines": ["不 会! 水 一 直 流",
                       "No — water keeps flowing",
                       "用 不 完 的 能 源 · Renewable!"]},
        ],
        "frame_cn": "「水 力 来 自 流 动 的 水.」",
        "frame_en": "Hydro energy comes from flowing water.",
    },
    {
        "title": "🛢️ 煤 / 石 油", "title_en": "Coal / oil",
        "img_label": "📸 一 堆 黑 煤 + 油 桶", "color": FOSSIL,
        "panels": [
            {"mark": "⚠️", "q": "这 是 什 么 能 源?",
             "lines": ["化 石 能 源 — 煤 和 石 油",
                       "Fossil energy — coal & oil",
                       "烧 了 才 有 能 量 · Burned for energy"]},
            {"mark": "❌", "q": "会 用 完 吗?",
             "lines": ["会! 挖 完 了 就 没 了",
                       "Yes — once dug up, it's gone",
                       "还 会 有 黑 烟 污 染 · Makes smoke"]},
        ],
        "frame_cn": "「化 石 能 源 会 用 完.」",
        "frame_en": "Fossil energy runs out.",
    },
    {
        "title": "🚗 太 阳 能 小 车", "title_en": "Solar toy car",
        "img_label": "📸 一 辆 太 阳 能 小 车", "color": SOLAR,
        "panels": [
            {"mark": "✅", "q": "这 是 什 么 能 源?",
             "lines": ["太 阳 能! 车 顶 有 太 阳 能 板",
                       "Solar — the panel is on top",
                       "晒 到 太 阳 就 会 跑 · Runs in sunlight"]},
            {"mark": "🌞", "q": "会 用 完 吗?",
             "lines": ["不 会! 只 要 有 太 阳",
                       "No — as long as there's sun",
                       "下 午 我 们 就 来 做 一 个! · We'll build one!"]},
        ],
        "frame_cn": "「太 阳 能 来 自 太 阳.」",
        "frame_en": "Solar energy comes from the sun.",
    },
]

VOCAB_RECOGNIZE = [
    ("🌍", "地 球", "dì qiú",     "Earth",
     "我 们 住 在 地 球 上.", "We live on Earth.",
     "📸 蓝 色 的 地 球", DAY),
    ("⛏️", "资 源", "zī yuán",   "resources",
     "地 球 的 资 源 有 限.", "Earth's resources are limited.",
     "📸 水 / 树 / 矿 石", EARTH_BROWN),
    ("⚡", "能 源", "néng yuán",  "energy",
     "灯 和 车 都 需 要 能 源.", "Lights and cars need energy.",
     "📸 闪 电 / 电 池 / 插 头", FIRE_ORANGE),
    ("☀️", "太 阳 能", "tài yáng néng", "solar energy",
     "太 阳 能 来 自 太 阳.", "Solar energy comes from the sun.",
     "📸 太 阳 + 太 阳 能 板", SOLAR),
    ("💡", "节 约", "jié yuē",    "save / conserve",
     "我 们 要 节 约 用 电.", "We should save electricity.",
     "📸 关 灯 的 手", MOSS),
]

VOCAB_WRITE = [
    ("地 球", "Earth", DAY, [
        ("地", "dì",  "6 笔 / 6 strokes", "土 字 旁 + 也 — 土 地 的 地"),
        ("球", "qiú", "11 笔 / 11 strokes", "王 (玉) 字 旁 + 求 — 圆 圆 的 球"),
    ]),
    ("节 约", "save", MOSS, [
        ("节", "jié", "5 笔 / 5 strokes", "艹 头 + 卩 — 像 竹 子 一 节 节"),
        ("约", "yuē", "6 笔 / 6 strokes", "纟 (绞 丝 旁) + 勺 — 约 束 一 下"),
    ]),
]


# ============================================================
# 1 · COVER
# ============================================================
cover(prs, 2, "能 量 从 哪 里 来?", "Where Does Energy Come From?",
      "☀️  💨  💧  🔋",
      DAY,
      "用 不 完 的 能 源 在 哪 里?",
      "Where can we find energy that never runs out?")
n += 1; pn(prs.slides[-1], n)
notes(prs.slides[-1], "📍 Day 2 · 能 量 从 哪 里 来? (可 再 生 能 源)\n👩‍🏫 老 师 说: 「今 天 我 们 是 小 小 能 源 工 程 师 — 帮 地 球 找 用 不 完 的 能 源!」\n⏱️ 1 分 钟")


# ============================================================
# 2 · SESSION 1 DIVIDER
# ============================================================
s = div(prs, "Session 1",
        "🔍 上 午 10:00–10:45 / 11:00–11:45  ·  能 量 从 哪 里 来?",
        DAY, "⚡"); n += 1; pn(s, n)


# ============================================================
# 3 · LEARNING GOALS
# ============================================================
s = learning_goals(prs, DAY, [
    ("1️⃣", "明 白 地 球 的 资 源 有 限 — 有 些 能 源 会 用 完",
     "Earth's resources are limited", FIRE_ORANGE),
    ("2️⃣", "初 步 懂 得 「可 持 续」 — 用 不 完 的 能 源 更 好",
     "Sustainability — pick energy that lasts", MOSS),
    ("3️⃣", "认 识 太 阳 能 · 风 能 · 水 力",
     "Know solar, wind & hydro energy", SOLAR),
    ("4️⃣", "学 会 句 型: 「___ 能 来 自 ___」「我 会 节 约 ___」",
     "Use the sentence frames", DAY),
])
n += 1; pn(s, n)


# ============================================================
# 4 · HOOK — 什么需要能量?
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "⚡ 热 身  ·  什 么 需 要 能 量?", DAY)
tb(s, 0.4, 0.85, 9.2, 0.28,
   "看 一 看 — 这 些 东 西 要 动 起 来, 需 要 什 么?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.15, 9.2, 0.24,
   "What do all these things need to work?",
   sz=10, c=GRAY, a=PP_ALIGN.CENTER)
hook_items = [
    ("💡", "电 灯", "Light", SOLAR),
    ("🚗", "汽 车", "Car", FIRE_ORANGE),
    ("📱", "手 机", "Phone", WIND),
    ("📺", "电 视", "TV", HYDRO),
]
hw = 2.10; hgap = 0.20
htotal = 4*hw + 3*hgap; hstart = (10 - htotal)/2
for i, (em, cn, en, cl) in enumerate(hook_items):
    x = hstart + i*(hw + hgap)
    panel(s, x, 1.55, hw, 1.95, cl, fill=WHITE, lw=2.5)
    tb(s, x, 1.70, hw, 0.80, em, sz=48, a=PP_ALIGN.CENTER)
    tb(s, x, 2.55, hw, 0.42, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.00, hw-0.10, 0.30, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
ans = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.80), Inches(9.20), Inches(1.55))
ans.fill.solid(); ans.fill.fore_color.rgb = DAY
ans.line.color.rgb = STAR; ans.line.width = Pt(3)
tb(s, 0.55, 3.95, 9.0, 0.55, "⚡ 都 需 要 — 能 量 (能 源)!",
   sz=26, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.55, 9.0, 0.32, "They all need ENERGY!",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.90, 9.0, 0.32,
   "🤔 那 — 能 量 从 哪 里 来 呢?  Where does energy come from?",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "热 身 3-4 分 钟\n👩‍🏫 问: 「这 些 东 西 不 动 了 — 少 了 什 么?」\n💡 引 出 「能 量 / 能 源」, 先 不 讲 来 源 — 留 给 下 面")


# ============================================================
# 5 · 能源是什么 (concept)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "💡 能 源 是 什 么?  ·  What Is Energy?", DAY)
tb(s, 0.4, 0.90, 9.2, 0.45,
   "能 源 = 让 东 西 「动 起 来」 的 力 量",
   sz=22, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.40, 9.2, 0.30,
   "Energy is the power that makes things work.",
   sz=12, c=GRAY, a=PP_ALIGN.CENTER)
concept = [
    ("🔋", "能 源 让 灯 亮", "Energy lights the lamp", SOLAR),
    ("🚗", "能 源 让 车 跑", "Energy moves the car", FIRE_ORANGE),
    ("🏭", "能 源 让 机 器 工 作", "Energy runs machines", HYDRO),
]
cw = 2.90; cgap = 0.15
ctotal = 3*cw + 2*cgap; cstart = (10 - ctotal)/2
for i, (em, cn, en, cl) in enumerate(concept):
    x = cstart + i*(cw + cgap)
    panel(s, x, 1.95, cw, 2.20, cl, fill=WHITE, lw=2.5)
    tb(s, x, 2.10, cw, 0.85, em, sz=52, a=PP_ALIGN.CENTER)
    tb(s, x, 3.00, cw, 0.45, cn, sz=16, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.55, cw-0.10, 0.40, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.45), Inches(9.20), Inches(0.95))
qb.fill.solid(); qb.fill.fore_color.rgb = INK
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 4.55, 9.0, 0.45, "🤔 这 么 多 能 量 — 都 从 哪 里 来?",
   sz=20, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.02, 9.0, 0.30, "Where does all this energy come from?",
   sz=11, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "概 念 2-3 分 钟\n💡 用 身 边 例 子: 教 室 的 灯, 校 车")


# ============================================================
# 6 · THINK-PAIR-SHARE
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🤝 想 一 想 · 说 一 说  ·  Think-Pair-Share", DAY)
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(0.95), Inches(9.20), Inches(1.05))
qb.fill.solid(); qb.fill.fore_color.rgb = DAY
qb.line.color.rgb = STAR; qb.line.width = Pt(3)
tb(s, 0.55, 1.08, 9.0, 0.45, "❓ 能 量 可 能 从 哪 里 来?",
   sz=24, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 1.55, 9.0, 0.35, "Where might energy come from?",
   sz=13, c=WARM, a=PP_ALIGN.CENTER)
steps = [
    ("🧠", "想 一 想", "Think · 自 己 想 30 秒", HYDRO),
    ("👥", "两 人 说", "Pair · 跟 同 桌 说 1 分 钟", MOSS),
    ("🎤", "全 班 听", "Share · 几 位 同 学 分 享", FIRE_ORANGE),
]
sw = 2.85; sgap = 0.20
stotal = 3*sw + 2*sgap; sstart = (10 - stotal)/2
for i, (em, cn, en, cl) in enumerate(steps):
    x = sstart + i*(sw + sgap)
    panel(s, x, 2.25, sw, 1.85, cl, fill=WHITE, lw=2.5)
    tb(s, x, 2.40, sw, 0.70, em, sz=44, a=PP_ALIGN.CENTER)
    tb(s, x, 3.15, sw, 0.38, cn, sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.58, sw-0.20, 0.42, en, sz=10, c=GRAY, a=PP_ALIGN.CENTER)
tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.05))
tbox.fill.solid(); tbox.fill.fore_color.rgb = STAR
tbox.line.color.rgb = DAY; tbox.line.width = Pt(3)
tb(s, 0.55, 4.42, 9.0, 0.40, "⏱️ 计 时: 30 秒 · 60 秒 · 60 秒",
   sz=16, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.85, 9.0, 0.35,
   "Timer: 30s think · 60s pair · 60s share  (~3 min)",
   sz=11, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Think-Pair-Share 3-4 分 钟\n💡 老 师 把 答 案 记 在 黑 板: 太 阳? 风? 电? 油? — 不 评 价")


# ============================================================
# 7 · 化石能源 zone (INK) — 煤·石油·天然气
# ============================================================
s = zone_slide(prs, "🛢️", "化 石 能 源", "Fossil Energy",
               FOSSIL,
               [("⛏️", "煤 — 从 地 下 挖 出 来", "Coal — dug from the ground"),
                ("🛢️", "石 油 — 从 地 下 抽 出 来", "Oil — pumped from the ground"),
                ("🔥", "天 然 气 — 烧 了 才 有 能 量", "Natural gas — burned for energy"),
                ("🚗", "汽 车 · 工 厂 都 在 用", "Cars & factories use it"),
                ("⚠️", "现 在 用 得 最 多 — 但 有 问 题…", "Most-used today — but there's a problem…")],
               "「化 石 能 源 来 自 地 下.」",
               "Fossil energy comes from underground.",
               img_label="📸 煤 矿 / 油 井 — 真 实 照 片")
n += 1; pn(s, n)
notes(s, "化 石 能 源 3 分 钟\n💡 重 点: 煤/石 油/天 然 气 = 现 在 用 得 最 多 的 能 源\n下 一 页 揭 出 「问 题」: 会 用 完 + 污 染")


# ============================================================
# 8 · 用得完 + 污染
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "⚠️ 化 石 能 源 的 2 个 问 题", FOSSIL)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "煤 和 石 油 很 好 用 — 但 是 有 大 问 题!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
probs = [
    ("⏳", "会 用 完!", "Runs out!",
     "挖 完 了 就 没 有 了 — 不 会 再 长 出 来.",
     "Once it's gone, it's gone forever.", FIRE_ORANGE),
    ("🏭", "会 污 染!", "Pollutes!",
     "烧 的 时 候 冒 黑 烟 — 弄 脏 空 气.",
     "Burning makes smoke that dirties the air.", FOSSIL),
]
pw = 4.45; pgap = 0.30
pstart = (10 - 2*pw - pgap)/2
for i, (em, cn, en, line_cn, line_en, cl) in enumerate(probs):
    x = pstart + i*(pw + pgap)
    panel(s, x, 1.30, pw, 3.30, cl, fill=WHITE, lw=3)
    tb(s, x, 1.50, pw, 1.05, em, sz=72, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, pw, 0.55, cn, sz=26, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.28, pw-0.10, 0.30, en, sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.20, 3.70, pw-0.40, 0.45, line_cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
    tb(s, x+0.20, 4.15, pw-0.40, 0.35, line_en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.75), Inches(9.20), Inches(0.65))
bb.fill.solid(); bb.fill.fore_color.rgb = INK
bb.line.color.rgb = STAR; bb.line.width = Pt(2.5)
tb(s, 0.55, 4.84, 9.0, 0.45,
   "🌍 地 球 的 资 源 有 限 — 我 们 要 找 用 不 完 的 能 源!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "2 个 问 题 3 分 钟\n💡 目 标 1: 地 球 资 源 有 限\n用 手 比 划: 「挖 — 挖 — 没 了!」\n引 出 下 一 页 的 对 比")


# ============================================================
# 9 · 用得完 vs 用不完 (compare_slide — NEW helper)
# ============================================================
s = compare_slide(prs, "⚖️ 用 得 完  vs  用 不 完  ·  Runs Out vs Never Runs Out",
    left={
        "tag_cn": "化 石 能 源", "tag_en": "Fossil",
        "badge": "❌ 会 用 完", "color": FOSSIL,
        "emoji": "🛢️ ⛏️ 🔥",
        "title_cn": "用 得 完", "title_en": "Runs Out",
        "bullets": [
            ("⏳", "挖 完 就 没 了", "Gets used up"),
            ("🏭", "烧 了 会 污 染", "Pollutes when burned"),
            ("🌍", "对 地 球 不 好", "Bad for Earth"),
        ],
    },
    right={
        "tag_cn": "可 再 生 能 源", "tag_en": "Renewable",
        "badge": "✅ 用 不 完", "color": MOSS,
        "emoji": "☀️ 💨 💧",
        "title_cn": "用 不 完", "title_en": "Never Runs Out",
        "bullets": [
            ("🔁", "太 阳/风/水 一 直 来", "Sun/wind/water keep coming"),
            ("🌱", "干 净 — 不 污 染", "Clean — no pollution"),
            ("💚", "对 地 球 好", "Good for Earth"),
        ],
    },
    frame_cn="「我 们 要 用 「用 不 完」 的 能 源.」",
    frame_en="We should use energy that never runs out.")
n += 1; pn(s, n)
notes(s, "对 比 4 分 钟 — 全 课 的 核 心!\n💡 目 标 1+2: 资 源 有 限 → 选 可 再 生\n左 边 (黑): 用 得 完 + 污 染; 右 边 (绿): 用 不 完 + 干 净\n问: 「你 想 用 哪 种? 为 什 么?」")


# ============================================================
# 10 · 可持续发展 (concept)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🌱 可 持 续  ·  Sustainability", DAY)
tb(s, 0.4, 0.90, 9.2, 0.45,
   "可 持 续 = 用 「用 不 完」 的 能 源, 地 球 一 直 都 够 用",
   sz=17, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.4, 1.40, 9.2, 0.30,
   "Sustainability = using energy that lasts, so Earth keeps enough for everyone.",
   sz=11, c=GRAY, a=PP_ALIGN.CENTER)
sust = [
    ("☀️", "太 阳 能", "Solar", SOLAR),
    ("💨", "风 能", "Wind", WIND),
    ("💧", "水 力", "Hydro", HYDRO),
]
uw = 2.90; ugap = 0.15
ustart = (10 - 3*uw - 2*ugap)/2
for i, (em, cn, en, cl) in enumerate(sust):
    x = ustart + i*(uw + ugap)
    panel(s, x, 1.95, uw, 2.05, cl, fill=WHITE, lw=3)
    tb(s, x, 2.10, uw, 0.85, em, sz=56, a=PP_ALIGN.CENTER)
    tb(s, x, 3.00, uw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.52, uw-0.10, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.30), Inches(9.20), Inches(1.05))
bb.fill.solid(); bb.fill.fore_color.rgb = MOSS
bb.line.color.rgb = STAR; bb.line.width = Pt(3)
tb(s, 0.55, 4.42, 9.0, 0.45, "🌍 这 3 种 都 「用 不 完」 — 我 们 来 一 个 一 个 认 识!",
   sz=15, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.92, 9.0, 0.32, "All 3 never run out — let's meet them one by one!",
   sz=11, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "可 持 续 2-3 分 钟\n💡 目 标 2: 不 用 把 「可 持 续 发 展」 讲 太 深 — 「用 不 完 = 地 球 一 直 够 用」 即 可")


# ============================================================
# 11-13 · 三种可再生能源 (zone_slide)
# ============================================================
for r in RENEWABLES:
    s = zone_slide(prs, r["emoji"], r["name_cn"], r["name_en"],
                   r["color"], r["examples"],
                   r["frame_cn"], r["frame_en"],
                   img_label=f"📸 {r['name_cn']} — 真 实 照 片")
    n += 1; pn(s, n)
    notes(s, f"{r['name_cn']} 2-3 分 钟\n💡 用 句 型 「{r['name_cn']} 来 自 ___」 跟 读 3 遍")


# ============================================================
# 14 · 猜一猜 A/B/C  (ab3_slide)
# ============================================================
s = ab3_slide(prs,
    "猜 一 猜!", "Guess the Energy!",
    "这 个 大 风 车 用 什 么 能 源 发 电?",
    "What energy does this turbine use?",
    [("A", "☀️", "太 阳 能", "Solar", SOLAR),
     ("B", "💨", "风 能", "Wind", WIND),
     ("C", "💧", "水 力", "Hydro", HYDRO)],
    color=DAY)
n += 1; pn(s, n)
notes(s, "猜 一 猜 2 分 钟\n💡 答 案: B 风 能! 风 吹 → 风 车 转 → 发 电\n让 学 生 用 句 型: 「风 能 来 自 风」")


# ============================================================
# 15-19 · 这是什么能源? ×5 (answer_panels_slide)
# ============================================================
for p in PRACTICE:
    s = answer_panels_slide(prs,
        f"{p['title']}  ·  {p['title_en']}",
        p["color"],
        p["panels"],
        img_label=p["img_label"],
        subtitle="想 一 想 → 看 答 案 — Think, then check the answer!",
        frame_cn=p["frame_cn"],
        frame_en=p["frame_en"])
    n += 1; pn(s, n)
    notes(s, f"{p['title']} 1-2 分 钟\n💡 全 班 用 句 型 说 一 遍")


# ============================================================
# 20 · 小组游戏 — 能源大分类
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🏆 小 组 游 戏  ·  能 源 大 分 类!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "把 8 张 卡 片 分 一 分 — 哪 些 「用 不 完」? 哪 些 「用 得 完」?",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
# Two bins
bins = [
    ("✅ 用 不 完", "Never runs out", MOSS, 0.60),
    ("❌ 用 得 完", "Runs out", FOSSIL, 5.10),
]
for cn, en, cl, x in bins:
    panel(s, x, 1.30, 4.30, 1.20, cl, fill=WHITE, lw=3)
    tb(s, x, 1.45, 4.30, 0.50, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.00, 4.20, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.4, 2.70, 9.2, 0.30, "🎴 卡 片 (每 组 一 套 · 8 张):",
   sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
cards = [
    ("☀️", "太 阳"), ("💨", "风"),
    ("💧", "流 水"), ("🌊", "瀑 布"),
    ("⛏️", "煤"), ("🛢️", "石 油"),
    ("🔥", "天 然 气"), ("🔆", "太 阳 能 板"),
]
cw = 1.05; cgap = 0.08
ctotal = 8*cw + 7*cgap; cstart = (10 - ctotal)/2
for i, (em, cn) in enumerate(cards):
    x = cstart + i*(cw + cgap)
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(3.10), Inches(cw), Inches(1.20))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = DAY; cd.line.width = Pt(1.5)
    tb(s, x, 3.20, cw, 0.55, em, sz=24, a=PP_ALIGN.CENTER)
    tb(s, x+0.02, 3.75, cw-0.04, 0.40, cn, sz=9, b=True, c=DAY, a=PP_ALIGN.CENTER)
inst = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.85))
inst.fill.solid(); inst.fill.fore_color.rgb = DAY
inst.line.color.rgb = STAR; inst.line.width = Pt(2.5)
tb(s, 0.55, 4.63, 9.0, 0.32,
   "👥 全 班 分 4-5 组 · 把 卡 片 放 到 对 的 一 边!",
   sz=12, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.97, 9.0, 0.30,
   "Split into groups · sort the 8 cards into the right bin",
   sz=9, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "分 类 游 戏 5-6 分 钟\n💡 答 案: 用 不 完 = 太 阳/风/流 水/瀑 布/太 阳 能 板; 用 得 完 = 煤/石 油/天 然 气")


# ============================================================
# 21 · 公布答案 + medals
# ============================================================
s = ns(prs); bg(s, INK); hb(s, "🏆 公 布 答 案 + 比 一 比!", FIRE_ORANGE)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "老 师 念 答 案 — 每 组 数 一 数 你 们 对 了 几 张!",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
ak = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(1.30), Inches(9.20), Inches(0.95))
ak.fill.solid(); ak.fill.fore_color.rgb = WHITE
ak.line.color.rgb = STAR; ak.line.width = Pt(2.5)
tb(s, 0.55, 1.38, 9.0, 0.28, "✅ 答 案  Answer Key:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 1.66, 9.0, 0.30,
   "✅ 用 不 完: ☀️太 阳 · 💨风 · 💧流 水 · 🌊瀑 布 · 🔆太 阳 能 板",
   sz=11, b=True, c=MOSS, a=PP_ALIGN.LEFT)
tb(s, 0.55, 1.94, 9.0, 0.30,
   "❌ 用 得 完: ⛏️煤 · 🛢️石 油 · 🔥天 然 气",
   sz=11, b=True, c=FIRE_ORANGE, a=PP_ALIGN.LEFT)
medals = [
    ("🥇", "第 一 名", "1st Place", GOLD_MEDAL),
    ("🥈", "第 二 名", "2nd Place", SILVER_MEDAL),
    ("🥉", "第 三 名", "3rd Place", BRONZE_MEDAL),
]
mw = 2.85; mgap = 0.20
mtotal = 3*mw + 2*mgap; mstart = (10 - mtotal)/2
for i, (em, cn, en, cl) in enumerate(medals):
    x = mstart + i*(mw + mgap)
    panel(s, x, 2.45, mw, 1.95, cl, fill=WHITE, lw=3)
    tb(s, x, 2.60, mw, 0.85, em, sz=58, a=PP_ALIGN.CENTER)
    tb(s, x, 3.45, mw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.95, mw-0.20, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.60), Inches(9.20), Inches(0.80))
cb.fill.solid(); cb.fill.fore_color.rgb = STAR; cb.line.fill.background()
tb(s, 0.55, 4.70, 9.0, 0.35,
   "🎉 🎊  全 班 鼓 掌!  你 们 都 是 小 小 能 源 工 程 师!  🎉",
   sz=14, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.08, 9.0, 0.22,
   "Big applause — every team is a Little Energy Engineer!",
   sz=9, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "公 布 答 案 3-4 分 钟\n💡 重 点 不 是 输 赢")


# ============================================================
# 22 · EXIT TICKET
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎫 出 门 票  ·  用 句 型 说 一 个!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "请 5-8 位 同 学 — 指 一 种 能 源, 用 句 型 说 一 句!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.30), Inches(8.80), Inches(1.55))
fr1.fill.solid(); fr1.fill.fore_color.rgb = DAY
fr1.line.color.rgb = STAR; fr1.line.width = Pt(3)
tb(s, 0.75, 1.42, 8.50, 0.35, "1️⃣  句 型 一  ·  Sentence Frame 1",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.85, 8.50, 0.60, "「______ 能 来 自 ______.」",
   sz=28, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.50, 8.50, 0.25, '"___ energy comes from ___."',
   sz=11, c=WARM, a=PP_ALIGN.CENTER)
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.10), Inches(9.20), Inches(2.25))
ex.fill.solid(); ex.fill.fore_color.rgb = WARM
ex.line.color.rgb = DAY; ex.line.width = Pt(2)
tb(s, 0.55, 3.20, 9.00, 0.30, "💡 举 个 例 子  Examples:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
exs = [
    "☀️ 「太 阳 能 来 自 太 阳.」",
    "💨 「风 能 来 自 风.」",
    "💧 「水 力 来 自 流 动 的 水.」",
]
for i, e in enumerate(exs):
    tb(s, 0.70, 3.55 + i*0.52, 8.70, 0.45, e, sz=16, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.08, 9.0, 0.25,
   "👏 说 对 了 — 你 就 是 今 天 的 小 小 能 源 工 程 师!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Exit Ticket 5-6 分 钟 (5-8 位 学 生)\n💡 K 级: 说 「太 阳 能」 即 可; G1-3: 整 句")


# ============================================================
# 23 · SESSION 2 DIVIDER
# ============================================================
s = div(prs, "Session 2",
        "📚 下 午 2:00–2:45  ·  词 汇 + 节 约 + 游 戏",
        DAY, "🔤"); n += 1; pn(s, n)


# ============================================================
# 24-28 · 我会认 × 5
# ============================================================
for em, cn, py, en, ex_cn, ex_en, hint, cl in VOCAB_RECOGNIZE:
    s = vocab_recognize(prs, cl, em, cn, py, en, ex_cn, ex_en, hint)
    n += 1; pn(s, n)
    notes(s, f"我 会 认 · {cn} 2-3 分 钟\n💡 跟 读 3 遍 + 造 句")


# ============================================================
# 29 · 词汇配对
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔗 配 对 游 戏  ·  Match the Words", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "把 字 跟 图 连 起 来!  Match each word with its picture!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
left_words = [
    ("地 球", DAY),
    ("资 源", EARTH_BROWN),
    ("能 源", FIRE_ORANGE),
    ("太 阳 能", SOLAR),
    ("节 约", MOSS),
]
right_pics = [
    ("💡", "C"), ("🌍", "E"), ("☀️", "B"), ("⛏️", "D"), ("⚡", "A"),
]
for i, (w, cl) in enumerate(left_words):
    y = 1.30 + i * 0.78
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.60), Inches(y), Inches(3.50), Inches(0.65))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = cl; box.line.width = Pt(2.5)
    tb(s, 0.60, y+0.10, 0.50, 0.45, str(i+1),
       sz=18, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, 1.15, y+0.10, 2.85, 0.45, w,
       sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
for i, (em, letter) in enumerate(right_pics):
    y = 1.30 + i * 0.78
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.90), Inches(y), Inches(3.50), Inches(0.65))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DAY; box.line.width = Pt(2.5)
    tb(s, 5.90, y+0.10, 0.50, 0.45, letter,
       sz=18, b=True, c=DAY, a=PP_ALIGN.CENTER)
    tb(s, 6.45, y+0.05, 2.85, 0.55, em, sz=28, a=PP_ALIGN.LEFT)
tb(s, 4.15, 3.10, 1.65, 0.40, "🤝 配 对",
   sz=14, b=True, c=DAY, a=PP_ALIGN.CENTER)
arr = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,
    Inches(4.15), Inches(3.50), Inches(1.65), Inches(0.40))
arr.fill.solid(); arr.fill.fore_color.rgb = STAR; arr.line.fill.background()
ak2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(5.20), Inches(9.20), Inches(0.30))
ak2.fill.solid(); ak2.fill.fore_color.rgb = DAY; ak2.line.fill.background()
tb(s, 0.55, 5.22, 9.0, 0.25,
   "🔑 答 案 (老 师 用): 1-E  2-D  3-A  4-B  5-C",
   sz=10, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "配 对 3-4 分 钟\n💡 1地球-E🌍 · 2资源-D⛏️ · 3能源-A⚡ · 4太阳能-B☀️ · 5节约-C💡")


# ============================================================
# 30 · 拍词卡
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "👋 拍 词 卡!  ·  Slap the Word!", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "两 人 一 组 — 老 师 念 词, 谁 先 拍 到 谁 赢!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
slap_cards = [
    ("🌍", "地 球", DAY),
    ("⛏️", "资 源", EARTH_BROWN),
    ("⚡", "能 源", FIRE_ORANGE),
    ("☀️", "太 阳 能", SOLAR),
    ("💡", "节 约", MOSS),
    ("👋", "拍!", STAR),
]
scw = 2.85; scgap = 0.15
scstart = (10 - 3*scw - 2*scgap) / 2
for i, (em, cn, cl) in enumerate(slap_cards):
    row = i // 3; col = i % 3
    x = scstart + col*(scw + scgap)
    y = 1.30 + row*1.40
    cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(scw), Inches(1.20))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE
    cd.line.color.rgb = cl; cd.line.width = Pt(3)
    tb(s, x+0.15, y+0.20, 0.90, 0.85, em, sz=44, a=PP_ALIGN.LEFT)
    tb(s, x+1.10, y+0.30, scw-1.20, 0.55, cn,
       sz=20, b=True, c=cl, a=PP_ALIGN.LEFT)
how_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.40), Inches(9.20), Inches(1.00))
how_box.fill.solid(); how_box.fill.fore_color.rgb = DAY
how_box.line.color.rgb = STAR; how_box.line.width = Pt(2.5)
tb(s, 0.55, 4.50, 9.00, 0.30, "🎯 怎 么 玩:",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.80, 9.00, 0.30,
   "• 两 人 一 组, 桌 上 摆 5 张 词 卡   • 老 师 念 词   • 谁 先 拍 中 谁 赢!",
   sz=12, b=True, c=WHITE, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.13, 9.00, 0.25,
   "Pairs · 5 cards on table · teacher calls · first slap wins",
   sz=9, c=WARM, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "拍 词 卡 4-5 分 钟\n💡 轻 拍! 不 是 打")


# ============================================================
# 31 · 句型练习 + Pair Share (both frames)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🗣️ 句 型 练 习  ·  Sentence Practice", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "和 同 桌 一 起 — 用 两 个 句 型 各 说 一 句!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
f1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.50), Inches(1.30), Inches(9.00), Inches(1.45))
f1.fill.solid(); f1.fill.fore_color.rgb = DAY
f1.line.color.rgb = STAR; f1.line.width = Pt(3)
tb(s, 0.65, 1.40, 8.70, 0.32, "1️⃣  认 识 能 源",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.65, 1.74, 8.70, 0.55, "「______ 能 来 自 ______.」",
   sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.65, 2.34, 8.70, 0.30, "☀️ 太 阳 能 来 自 太 阳.  /  💨 风 能 来 自 风.",
   sz=12, b=True, c=WARM, a=PP_ALIGN.CENTER)
f2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.50), Inches(2.90), Inches(9.00), Inches(1.45))
f2.fill.solid(); f2.fill.fore_color.rgb = MOSS
f2.line.color.rgb = STAR; f2.line.width = Pt(3)
tb(s, 0.65, 3.00, 8.70, 0.32, "2️⃣  节 约 资 源",
   sz=12, b=True, c=STAR, a=PP_ALIGN.LEFT)
tb(s, 0.65, 3.34, 8.70, 0.55, "「我 会 节 约 ______.」",
   sz=24, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.65, 3.94, 8.70, 0.30, "💡 我 会 节 约 电.  /  💧 我 会 节 约 水.",
   sz=12, b=True, c=WARM, a=PP_ALIGN.CENTER)
tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.50), Inches(9.20), Inches(0.90))
tip.fill.solid(); tip.fill.fore_color.rgb = WARM
tip.line.color.rgb = DAY; tip.line.width = Pt(2)
tb(s, 0.55, 4.58, 9.0, 0.30, "👥 同 桌 练 习  Pair Practice:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
tb(s, 0.55, 4.88, 9.0, 0.42,
   "一 人 说 句 型 1, 一 人 说 句 型 2 — 然 后 交 换!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "句 型 练 习 5 分 钟 — 两 个 目 标 句 型\n💡 走 动 听, 请 2-3 对 上 来 示 范")


# ============================================================
# 32-33 · 我会写 × 2
# ============================================================
for cn_phrase, en_word, cl, chars in VOCAB_WRITE:
    s = vocab_write(prs, cl, cn_phrase, en_word, chars)
    n += 1; pn(s, n)
    notes(s, f"我 会 写 · {cn_phrase} 3-4 分 钟\n💡 看 老 师 写 → 空 中 写 → 田 字 格 写 3 遍")


# ============================================================
# 34 · 如何节约资源
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "💚 如 何 节 约 资 源?  ·  How to Save", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "能 源 很 宝 贵 — 我 们 每 天 都 可 以 节 约!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
saves = [
    ("💡", "关 灯", "Turn off lights", "走 出 房 间 就 关 灯", SOLAR),
    ("🚿", "关 水", "Turn off water", "刷 牙 时 关 水", HYDRO),
    ("♻️", "重 复 用", "Reuse", "纸 两 面 用 / 带 水 壶", MOSS),
]
vw = 2.90; vgap = 0.15
vstart = (10 - 3*vw - 2*vgap)/2
for i, (em, cn, en, line_cn, cl) in enumerate(saves):
    x = vstart + i*(vw + vgap)
    panel(s, x, 1.45, vw, 2.55, cl, fill=WHITE, lw=3)
    tb(s, x, 1.60, vw, 0.85, em, sz=56, a=PP_ALIGN.CENTER)
    tb(s, x, 2.50, vw, 0.45, cn, sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.02, vw-0.10, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.42, vw-0.20, 0.50, line_cn, sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
fb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.25), Inches(9.20), Inches(1.10))
fb.fill.solid(); fb.fill.fore_color.rgb = INK
fb.line.color.rgb = STAR; fb.line.width = Pt(3)
tb(s, 0.55, 4.36, 9.0, 0.32, "💬 句 型: 「我 会 节 约 ______.」",
   sz=15, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.78, 9.0, 0.45,
   "我 会 节 约 电.  ·  我 会 节 约 水.  ·  我 会 节 约 纸.",
   sz=15, b=True, c=WHITE, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "节 约 4 分 钟 — 目 标 4\n💡 每 个 学 生 说 一 句 「我 会 节 约 ___」")


# ============================================================
# 35 · BAMBOOZLE
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎮 复 习 游 戏  ·  Bamboozle", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "全 班 分 组 抢 答 — 复 习 今 天 的 能 源 知 识!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
panel(s, 0.40, 1.25, 4.40, 3.55, DAY, fill=INK, lw=3)
tb(s, 0.40, 1.55, 4.40, 1.10, "🎮", sz=80, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.75, 4.40, 0.45, "Bamboozle", sz=22, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.25, 4.40, 0.30, "能 源 大 复 习", sz=14, b=True, c=WARM, a=PP_ALIGN.CENTER)
pb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.20), Inches(3.80), Inches(2.80), Inches(0.60))
pb.fill.solid(); pb.fill.fore_color.rgb = FIRE_ORANGE; pb.line.fill.background()
tb(s, 1.20, 3.90, 2.80, 0.42, "▶️  点 击 打 开  Open",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
pb.click_action.hyperlink.address = "https://www.baamboozle.com/"
panel(s, 5.10, 1.25, 4.50, 3.55, FIRE_ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.10, 1.25, 4.50, FIRE_ORANGE, "👩‍🏫 老 师 准 备", sz=13)
setup = [
    "1. 打 开 baamboozle.com",
    "2. import 「bamboozle_day2_energy.csv」",
    "   (在 同 一 个 文 件 夹)",
    "3. 全 班 分 2-4 组",
    "4. 答 对 加 分, 答 错 一 起 复 习",
]
for i, t in enumerate(setup):
    tb(s, 5.25, 1.85 + i*0.52, 4.25, 0.45, t,
       sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.45))
nb.fill.solid(); nb.fill.fore_color.rgb = STAR; nb.line.fill.background()
tb(s, 0.55, 5.02, 9.0, 0.32,
   "🎯 10 题: 用 不 完/用 得 完 · 太 阳/风/水 · 节 约",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Bamboozle 8-10 分 钟\n💡 CSV 在 同 一 文 件 夹: bamboozle_day2_energy.csv")


# ============================================================
# 36 · CLOSING REFLECTION (Session 2)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "💭 今 天 我 学 会 了 ...  Today I Learned ...", DAY)
tb(s, 0.4, 0.90, 9.2, 0.40,
   "想 一 想 — 用 一 句 话 说 出 你 学 到 的!",
   sz=14, b=True, c=DARK, a=PP_ALIGN.CENTER)
rb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.55), Inches(8.80), Inches(1.50))
rb.fill.solid(); rb.fill.fore_color.rgb = DAY
rb.line.color.rgb = STAR; rb.line.width = Pt(3)
tb(s, 0.75, 1.95, 8.50, 0.70, "「今 天 我 学 会 了 ____________.」",
   sz=26, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.65, 8.50, 0.30, '"Today I learned ___."',
   sz=12, c=WARM, a=PP_ALIGN.CENTER)
ideas = [
    ("☀️", "太 阳 能 来 自 太 阳", SOLAR),
    ("💨", "风 能 来 自 风", WIND),
    ("💧", "水 力 来 自 水", HYDRO),
    ("💚", "我 会 节 约 资 源", MOSS),
]
iw = 2.20; igap = 0.15
istart = (10 - 4*iw - 3*igap)/2
for i, (em, cn, cl) in enumerate(ideas):
    x = istart + i*(iw + igap)
    panel(s, x, 3.30, iw, 1.65, cl, fill=WHITE, lw=2.5)
    tb(s, x, 3.45, iw, 0.65, em, sz=40, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 4.15, iw-0.10, 0.70, cn, sz=12, b=True, c=cl, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "反 思 3-4 分 钟\n💡 每 个 学 生 说 一 句, 老 师 不 评 价")


# ============================================================
# 37 · SESSION 3 DIVIDER
# ============================================================
s = div(prs, "Session 3",
        "🔧 下 午  ·  我 是 小 小 能 源 工 程 师!",
        DAY, "🔧"); n += 1; pn(s, n)


# ============================================================
# 38 · PROJECT MENU — 太阳能小车 / 太阳能发电
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🛠️ 今 天 做 什 么?  ·  Solar Project", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "我 们 来 做 一 个 太 阳 能 作 品! (老 师 看 套 件 选 一 个)",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
projs = [
    ("🚗", "太 阳 能 小 车", "Solar Car",
     "晒 太 阳 → 小 车 自 己 跑!", SOLAR,
     os.path.join(CAR_DIR, "2-效果图.png")),
    ("💡", "太 阳 能 发 电", "Solar Power",
     "晒 太 阳 → 灯 亮 / 风 扇 转!", FIRE_ORANGE,
     os.path.join(GEN_DIR, "2-效果图.png")),
]
pw = 4.45; pgap = 0.30
pstart = (10 - 2*pw - pgap)/2
for i, (em, cn, en, line_cn, cl, img) in enumerate(projs):
    x = pstart + i*(pw + pgap)
    panel(s, x, 1.30, pw, 3.95, cl, fill=WHITE, lw=3)
    tb(s, x, 1.42, pw, 0.55, f"{em} {cn}", sz=22, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 1.98, pw-0.10, 0.28, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    img_or_slot(s, x+0.30, 2.35, pw-0.60, 2.10, img,
                f"{cn} 成 品", en, cl)
    tb(s, x+0.15, 4.55, pw-0.30, 0.55, line_cn, sz=13, b=True, c=cl, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "选 项 1-2 分 钟\n💡 老 师 根 据 现 有 太 阳 能 套 件 二 选 一\n两 个 项 目 的 步 骤 图 / 视 频 都 在 单 元 文 件 夹 里")


# ============================================================
# 39 · MATERIALS (reuse 1-材料准备 images)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎒 材 料  ·  Materials", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "看 看 我 们 需 要 什 么 — 老 师 已 经 准 备 好 啦!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
mat = [
    ("🚗 太 阳 能 小 车", SOLAR, os.path.join(CAR_DIR, "1-彩色版", "1-材料准备.png")),
    ("💡 太 阳 能 发 电", FIRE_ORANGE, os.path.join(GEN_DIR, "1-彩色版", "1-材料准备.png")),
]
mw = 4.45; mgap = 0.30
mstart = (10 - 2*mw - mgap)/2
for i, (cn, cl, img) in enumerate(mat):
    x = mstart + i*(mw + mgap)
    panel(s, x, 1.25, mw, 3.95, cl, fill=WHITE, lw=3)
    tb(s, x, 1.35, mw, 0.45, cn, sz=16, b=True, c=cl, a=PP_ALIGN.CENTER)
    img_or_slot(s, x+0.25, 1.90, mw-0.50, 3.05, img,
                "材 料 准 备", "Materials", cl)
n += 1; pn(s, n)
notes(s, "材 料 2 分 钟\n💡 太 阳 能 套 件: 太 阳 能 板 + 马 达 + (小 车 轮 子 / LED 灯)\n图 来 自 单 元 「步 骤 图/1-彩 色 版/1-材 料 准 备」")


# ============================================================
# 40 · STEPS (reuse 步骤图 images)
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🔧 动 手 步 骤  ·  Build Steps", DAY)
tb(s, 0.4, 0.82, 9.2, 0.28,
   "跟 着 图 一 步 一 步 做 — 老 师 先 示 范!",
   sz=12, b=True, c=DARK, a=PP_ALIGN.CENTER)
# Car steps row
tb(s, 0.40, 1.12, 4.30, 0.30, "🚗 太 阳 能 小 车", sz=12, b=True, c=SOLAR, a=PP_ALIGN.LEFT)
car_steps = [
    os.path.join(CAR_DIR, "1-彩色版", "2-步骤1、2.png"),
    os.path.join(CAR_DIR, "1-彩色版", "3-步骤3、4.png"),
]
for i, img in enumerate(car_steps):
    img_or_slot(s, 0.40 + i*2.20, 1.45, 2.05, 1.75, img,
                f"小 车 步 骤 {i+1}", "Car step", SOLAR)
# Gen steps row
tb(s, 5.05, 1.12, 4.30, 0.30, "💡 太 阳 能 发 电", sz=12, b=True, c=FIRE_ORANGE, a=PP_ALIGN.LEFT)
gen_steps = [
    os.path.join(GEN_DIR, "1-彩色版", "2-步骤图1.png"),
    os.path.join(GEN_DIR, "1-彩色版", "3-步骤图2.png"),
]
for i, img in enumerate(gen_steps):
    img_or_slot(s, 5.05 + i*2.20, 1.45, 2.05, 1.75, img,
                f"发 电 步 骤 {i+1}", "Power step", FIRE_ORANGE)
# Bottom tips
tipb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.45), Inches(9.20), Inches(1.95))
tipb.fill.solid(); tipb.fill.fore_color.rgb = WHITE
tipb.line.color.rgb = DAY; tipb.line.width = Pt(2.5)
tb(s, 0.55, 3.55, 9.0, 0.35, "🔑 4 个 关 键 步 骤  4 Key Steps:",
   sz=13, b=True, c=DAY, a=PP_ALIGN.LEFT)
ksteps = [
    "1️⃣ 装 太 阳 能 板  Attach the solar panel",
    "2️⃣ 接 上 马 达 / 灯  Connect the motor / light",
    "3️⃣ 装 好 零 件  Assemble the parts",
    "4️⃣ 拿 到 太 阳 下 — 试 一 试!  Take it to the sun — test!",
]
for i, t in enumerate(ksteps):
    tb(s, 0.70, 3.95 + i*0.36, 8.70, 0.32, t, sz=12, b=True, c=DARK, a=PP_ALIGN.LEFT)
n += 1; pn(s, n)
notes(s, "步 骤 4-5 分 钟\n💡 图 来 自 单 元 步 骤 图 (彩 色 版)\n小 车: 步 骤 1、2 / 3、4 · 发 电: 步 骤 图 1 / 2")


# ============================================================
# 41 · TEACHER DEMO — video
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "📺 老 师 示 范  ·  Teacher Demo Video", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "看 操 作 视 频 — 老 师 一 步 一 步 做 给 你 看!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
vids = [
    ("🚗 太 阳 能 小 车 操 作 视 频", SOLAR, CAR_VIDEO,
     "互动版_太阳能小车_科学区_大班/太阳能小车_科学区_大班_操作视频.mp4"),
    ("💡 太 阳 能 发 电 操 作 视 频", FIRE_ORANGE, GEN_VIDEO,
     "互动版_太阳能发电_科学区_大班/太阳能发电_科学区_大班_操作视频.mp4"),
]
vw = 4.45; vgap = 0.30
vstart = (10 - 2*vw - vgap)/2
for i, (cn, cl, vpath, rel) in enumerate(vids):
    x = vstart + i*(vw + vgap)
    panel(s, x, 1.30, vw, 3.40, cl, fill=INK, lw=3)
    tb(s, x, 1.55, vw, 1.05, "📺", sz=72, a=PP_ALIGN.CENTER)
    tb(s, x, 2.70, vw, 0.45, cn, sz=15, b=True, c=STAR, a=PP_ALIGN.CENTER)
    btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x+0.85), Inches(3.30), Inches(vw-1.70), Inches(0.55))
    btn.fill.solid(); btn.fill.fore_color.rgb = FIRE_ORANGE; btn.line.fill.background()
    tb(s, x+0.85, 3.40, vw-1.70, 0.40, "▶️ 点 击 播 放  Play",
       sz=13, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    try:
        btn.click_action.hyperlink.address = rel
    except Exception:
        pass
    tb(s, x+0.10, 4.00, vw-0.20, 0.55, rel, sz=8, c=LGRAY, a=PP_ALIGN.CENTER)
nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.85), Inches(9.20), Inches(0.55))
nb.fill.solid(); nb.fill.fore_color.rgb = STAR; nb.line.fill.background()
tb(s, 0.55, 4.93, 9.0, 0.38,
   "🎬 视 频 在 单 元 文 件 夹 — 老 师 打 开 对 应 项 目 的 操 作 视 频",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "示 范 3-4 分 钟\n💡 视 频 是 本 地 文 件 (单 元 文 件 夹). 若 链 接 打 不 开, 老 师 直 接 双 击 mp4 播 放")


# ============================================================
# 42 · WORK TIME
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "⏱️ 动 手 时 间!  ·  Work Time", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "现 在 你 来 做! 老 师 走 动 帮 忙.",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
panel(s, 0.40, 1.20, 4.60, 3.30, DAY, fill=WHITE, lw=3)
tb(s, 0.40, 1.40, 4.60, 1.50, "⏰", sz=110, a=PP_ALIGN.CENTER)
tb(s, 0.40, 2.95, 4.60, 0.55, "35 分 钟", sz=42, b=True, c=DAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.55, 4.60, 0.35, "35 minutes · Build time", sz=13, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, 0.40, 3.95, 4.60, 0.32, "🎵 老 师 放 一 首 背 景 音 乐", sz=10, b=True, c=DAY, a=PP_ALIGN.CENTER)
panel(s, 5.20, 1.20, 4.40, 3.30, FIRE_ORANGE, fill=WHITE, lw=3)
panel_head(s, 5.20, 1.20, 4.40, FIRE_ORANGE, "🚶 老 师 走 动 提 醒", sz=13)
walk_tips = [
    "✦ 太 阳 能 板 要 对 准 光",
    "✦ 接 线 接 牢 (老 师 帮 忙)",
    "✦ K-2 — 帮 忙 装 零 件",
    "✦ 3-5 — 说 「太 阳 能 来 自 太 阳」",
    "✦ 做 完 拿 到 窗 边/灯 下 试!",
]
for i, t in enumerate(walk_tips):
    tb(s, 5.35, 1.85 + i*0.50, 4.10, 0.40, t, sz=11, b=True, c=DARK, a=PP_ALIGN.LEFT)
rc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.65), Inches(9.20), Inches(0.75))
rc.fill.solid(); rc.fill.fore_color.rgb = STAR; rc.line.fill.background()
tb(s, 0.55, 4.75, 9.0, 0.32,
   "✅ 完 成 检 查: 太 阳 能 板 装 好 · 零 件 接 牢 · 能 在 光 下 动 起 来",
   sz=12, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.10, 9.0, 0.22,
   "Done check: panel attached · parts connected · works in the light",
   sz=9, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "动 手 35 分 钟\n💡 K-2 需 要 老 师 帮 接 线; 接 线 是 唯 一 需 要 帮 忙 的 步 骤")


# ============================================================
# 43 · 测试! 拿到阳光下
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "☀️ 测 试 时 间!  ·  Test in the Sun!", SOLAR)
tb(s, 0.4, 0.90, 9.2, 0.40,
   "拿 到 太 阳 下 (或 大 灯 下) — 看 它 动 起 来!",
   sz=16, b=True, c=DARK, a=PP_ALIGN.CENTER)
test = [
    ("☀️", "找 到 光", "Find the light", "走 到 窗 边 / 院 子 / 大 灯 下", SOLAR),
    ("🔆", "对 准 板", "Aim the panel", "太 阳 能 板 对 准 光", FIRE_ORANGE),
    ("🎉", "动 起 来!", "It works!", "小 车 跑 / 灯 亮 / 风 扇 转!", MOSS),
]
tw = 2.90; tgap = 0.15
tstart = (10 - 3*tw - 2*tgap)/2
for i, (em, cn, en, line_cn, cl) in enumerate(test):
    x = tstart + i*(tw + tgap)
    panel(s, x, 1.55, tw, 2.55, cl, fill=WHITE, lw=3)
    tb(s, x, 1.70, tw, 0.85, em, sz=56, a=PP_ALIGN.CENTER)
    tb(s, x, 2.60, tw, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.10, tw-0.10, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.48, tw-0.20, 0.50, line_cn, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.35), Inches(9.20), Inches(1.00))
bb.fill.solid(); bb.fill.fore_color.rgb = SOLAR
bb.line.color.rgb = INK; bb.line.width = Pt(2.5)
tb(s, 0.55, 4.45, 9.0, 0.40, "🤔 没 有 光 会 怎 样? 试 试 用 手 遮 住 板 子!",
   sz=14, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.90, 9.0, 0.32,
   "No light = no power! That's why we need the sun. 🌞",
   sz=11, b=True, c=INK, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "测 试 5 分 钟 — 最 兴 奋 的 环 节!\n💡 遮 光 实 验: 让 学 生 亲 眼 看 到 「有 太 阳 才 有 能 量」")


# ============================================================
# 44 · GALLERY WALK
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🖼️ 作 品 展 示  ·  Gallery Walk", DAY)
tb(s, 0.4, 0.90, 9.2, 0.40,
   "把 作 品 摆 在 桌 上 — 我 们 一 起 看 一 看!",
   sz=15, b=True, c=DARK, a=PP_ALIGN.CENTER)
gw = [
    ("🪑", "摆 桌", "Set up", "把 作 品 摆 在 桌 上", SOLAR),
    ("🚶", "静 走", "Quiet walk", "安 静 地 看 别 人 的", HYDRO),
    ("👏", "拍 桌", "Applaud", "看 到 喜 欢 的 轻 拍 桌 子", MOSS),
]
gww = 2.90; gwgap = 0.15
gwstart = (10 - 3*gww - 2*gwgap)/2
for i, (em, cn, en, line_cn, cl) in enumerate(gw):
    x = gwstart + i*(gww + gwgap)
    panel(s, x, 1.55, gww, 2.55, cl, fill=WHITE, lw=3)
    tb(s, x, 1.70, gww, 0.85, em, sz=56, a=PP_ALIGN.CENTER)
    tb(s, x, 2.60, gww, 0.45, cn, sz=20, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.10, gww-0.10, 0.30, en, sz=11, c=GRAY, a=PP_ALIGN.CENTER)
    tb(s, x+0.10, 3.48, gww-0.20, 0.50, line_cn, sz=11, b=True, c=DARK, a=PP_ALIGN.CENTER)
nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.35), Inches(9.20), Inches(1.00))
nb.fill.solid(); nb.fill.fore_color.rgb = DAY
nb.line.color.rgb = STAR; nb.line.width = Pt(3)
tb(s, 0.55, 4.55, 9.0, 0.45, "📸 老 师 拍 照 — 留 下 大 家 的 太 阳 能 作 品!",
   sz=14, b=True, c=STAR, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "Gallery Walk 5 分 钟\n💡 规 则: 静 走, 不 碰 别 人 的 作 品")


# ============================================================
# 45 · 分享句型
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🎤 分 享  ·  Share Your Work", DAY)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "请 几 位 同 学 — 拿 着 作 品, 用 句 型 说 一 说!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.60), Inches(1.35), Inches(8.80), Inches(1.55))
fr.fill.solid(); fr.fill.fore_color.rgb = SOLAR
fr.line.color.rgb = INK; fr.line.width = Pt(3)
tb(s, 0.75, 1.48, 8.50, 0.35, "💬 分 享 句 型  Share Frame",
   sz=12, b=True, c=INK, a=PP_ALIGN.LEFT)
tb(s, 0.75, 1.88, 8.50, 0.60, "「我 用 太 阳 能 ____________.」",
   sz=26, b=True, c=INK, a=PP_ALIGN.CENTER)
tb(s, 0.75, 2.52, 8.50, 0.28, '"I used solar energy to ___."',
   sz=11, c=DARK, a=PP_ALIGN.CENTER)
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(3.10), Inches(9.20), Inches(2.25))
ex.fill.solid(); ex.fill.fore_color.rgb = WARM
ex.line.color.rgb = DAY; ex.line.width = Pt(2)
tb(s, 0.55, 3.20, 9.00, 0.30, "💡 举 个 例 子  Examples:",
   sz=11, b=True, c=DAY, a=PP_ALIGN.LEFT)
exs = [
    "🚗 「我 用 太 阳 能 让 小 车 跑.」",
    "💡 「我 用 太 阳 能 让 灯 亮.」",
    "🌞 「太 阳 能 来 自 太 阳 — 用 不 完!」",
]
for i, e in enumerate(exs):
    tb(s, 0.70, 3.55 + i*0.52, 8.70, 0.45, e, sz=16, b=True, c=DARK, a=PP_ALIGN.LEFT)
tb(s, 0.55, 5.08, 9.0, 0.25,
   "👏 每 个 工 程 师 都 说 一 句!  Every engineer shares one sentence!",
   sz=11, b=True, c=DAY, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "分 享 5-6 分 钟\n💡 K 级: 「太 阳 能!」 + 指 作 品; G1-3: 整 句")


# ============================================================
# 46 · 在家节约 / 活动延伸
# ============================================================
s = ns(prs); bg(s, CREAM); hb(s, "🏠 在 家 也 能 做  ·  At Home", MOSS)
tb(s, 0.4, 0.85, 9.2, 0.30,
   "今 天 学 的, 回 家 也 可 以 做 — 你 也 是 小 小 能 源 工 程 师!",
   sz=13, b=True, c=DARK, a=PP_ALIGN.CENTER)
home = [
    ("💡", "随 手 关 灯", "Turn off lights", SOLAR),
    ("🚿", "节 约 用 水", "Save water", HYDRO),
    ("🔍", "找 找 太 阳 能", "Spot solar around you", FIRE_ORANGE),
    ("👨‍👩‍👧", "告 诉 家 人", "Tell your family", MOSS),
]
hw = 2.20; hgap = 0.15
hstart = (10 - 4*hw - 3*hgap)/2
for i, (em, cn, en, cl) in enumerate(home):
    x = hstart + i*(hw + hgap)
    panel(s, x, 1.45, hw, 2.40, cl, fill=WHITE, lw=2.5)
    tb(s, x, 1.60, hw, 0.85, em, sz=48, a=PP_ALIGN.CENTER)
    tb(s, x, 2.50, hw, 0.70, cn, sz=15, b=True, c=cl, a=PP_ALIGN.CENTER)
    tb(s, x+0.05, 3.25, hw-0.10, 0.45, en, sz=9, c=GRAY, a=PP_ALIGN.CENTER)
bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.40), Inches(4.10), Inches(9.20), Inches(1.25))
bb.fill.solid(); bb.fill.fore_color.rgb = INK
bb.line.color.rgb = STAR; bb.line.width = Pt(3)
tb(s, 0.55, 4.22, 9.0, 0.35, "🌍 小 任 务  Home Mission:",
   sz=13, b=True, c=STAR, a=PP_ALIGN.CENTER)
tb(s, 0.55, 4.60, 9.0, 0.40,
   "回 家 找 1 个 太 阳 能 的 东 西, 明 天 告 诉 大 家!",
   sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 0.55, 5.02, 9.0, 0.28,
   "Find 1 solar-powered thing at home — share tomorrow!",
   sz=10, c=WARM, a=PP_ALIGN.CENTER)
n += 1; pn(s, n)
notes(s, "活 动 延 伸 2-3 分 钟\n💡 家 园 共 育: 节 约 用 电 用 水 + 找 太 阳 能")


# ============================================================
# 47 · SHARE + CLOSE
# ============================================================
s = share_close(prs, DAY,
    ["「___ 能 来 自 ___.」", "「我 会 节 约 ___.」"],
    "___ energy comes from ___.  ·  I will save ___.",
    "更 多 爱 护 地 球 的 行 动!",
    "More ways to protect our Earth!",
    next_emoji="🌍")
n += 1; pn(s, n)
notes(s, "收 尾 2 分 钟\n💡 给 每 人 一 张 「小 小 能 源 工 程 师」 贴 纸 / 印 章\n预 告 下 一 天")


# ============================================================
# SAVE
# ============================================================
out = os.path.join(HERE, "day2_energy.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
