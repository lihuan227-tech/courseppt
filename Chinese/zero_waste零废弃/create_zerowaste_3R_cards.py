#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零废弃与 3R — 可打印教具  Printable Classroom Cards
Companion handout to zerowaste_3R.pptx.

Produces (US-Letter, cut on the dashed lines):
  Page 1     — 3 张 R 区标牌 (Reduce / Reuse / Recycle zone signs, 1 per row)
  Page 2-3   — 3R 分类图片卡 × 12 (sorting picture cards, 6 / page)
  Page 4-5   — Zero Waste 大侦探场景卡 × 8 (detective scene cards, 4 / page)
  Page 6     — 侦探记录单 (detective record sheet — 问题 / 更好的做法)
  Page 7     — 老师参考答案 (answer key)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# US Letter, portrait
prs.slide_width  = Inches(8.5)
prs.slide_height = Inches(11)
W, H = prs.slide_width, prs.slide_height

# --- shared palette with the deck ---
ECO     = RGBColor(0x2E,0x7D,0x32)
DEEP    = RGBColor(0x1B,0x5E,0x20)
REDUCE  = RGBColor(0x1E,0x88,0xE5)
REUSE   = RGBColor(0xFB,0x8C,0x00)
RECYCLE = RGBColor(0x43,0xA0,0x47)
BROWN   = RGBColor(0x6B,0x44,0x23)
ALERT   = RGBColor(0xD0,0x4A,0x3C)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x2C,0x2C,0x2C)
GRAY    = RGBColor(0x88,0x88,0x88)
LGRAY   = RGBColor(0xBB,0xBB,0xBB)
WARM    = RGBColor(0xF1,0xF8,0xE9)
CREAM   = RGBColor(0xFD,0xFB,0xF3)
GOLD    = RGBColor(0xF9,0xA8,0x25)

BASE = "/Users/Huan/0 projects/summercourse/Chinese/zero_waste零废弃"

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None,font='KaiTi'):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
    if a:p.alignment=a
    r=p.add_run();r.text=txt;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=font;return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    sp=sh._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp)
def dashed_rect(s,l,t,w,h,c=LGRAY,fill=WHITE,wd=1.25):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=fill
    sh.line.color.rgb=c;sh.line.width=Pt(wd);sh.line.dash_style=2
    return sh
def page_title(s,txt,sub,c=ECO):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.4),Inches(0.35),Inches(7.7),Inches(0.7))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    tb(s,0.55,0.45,7.4,0.5,txt,sz=22,b=True,c=WHITE)
    tb(s,0.4,1.10,7.7,0.3,sub,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
def scissor(s,t):
    tb(s,0.2,t-0.14,8.1,0.28,"✂ - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -",sz=11,c=LGRAY)

# ============================================================
# PAGE 1 — 3 R-zone signs (1 column, 3 rows)
# ============================================================
s=ns();bg(s,CREAM)
page_title(s,"🏷️ 3R 区标牌  Sorting Zone Signs","打印 → 剪开 → 贴在墙上或桌上 3 个区  ·  Print, cut, post 3 zones")
zones=[("⬇️","Reduce","减少","用得少一点  Use less",REDUCE),
       ("🔁","Reuse","重复使用","再用一次  Use again",REUSE),
       ("♻️","Recycle","回收","变成新东西  Make new",RECYCLE)]
top=1.6
for i,(em,en,cn,d,cl) in enumerate(zones):
    y=top+i*2.95
    scissor(s,y-0.12)
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(y),Inches(7.4),Inches(2.65))
    card.fill.solid();card.fill.fore_color.rgb=WHITE;card.line.color.rgb=cl;card.line.width=Pt(3)
    band=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(y),Inches(7.4),Inches(0.85))
    band.fill.solid();band.fill.fore_color.rgb=cl;band.line.fill.background()
    tb(s,0.7,y+0.12,7.1,0.6,f"{em}  {en} · {cn}",sz=30,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    tb(s,0.7,y+1.05,7.1,0.9,em,sz=64,a=PP_ALIGN.CENTER)
    tb(s,0.7,y+2.05,7.1,0.45,d,sz=20,b=True,c=cl,a=PP_ALIGN.CENTER)

# ============================================================
# PAGES 2-3 — 3R sorting picture cards (6 per page)
# ============================================================
# (emoji, 中文, English, teacher-answer color for the corner dot key on answer page)
sort_cards=[
    ("💧","一直流的水龙头","Running tap",REDUCE),
    ("📄","只写一面的纸","Paper, one side",REDUCE),
    ("🛍️","塑料袋","Plastic bag",REDUCE),
    ("🥤","塑料吸管","Plastic straw",REDUCE),
    ("🍶","空塑料瓶","Empty bottle",REUSE),
    ("👕","小了的衣服","Outgrown clothes",REUSE),
    ("📦","快递纸箱","Cardboard box",REUSE),
    ("🫙","玻璃罐子","Glass jar",REUSE),
    ("🥫","喝完的易拉罐","Empty can",RECYCLE),
    ("📰","旧报纸","Old newspaper",RECYCLE),
    ("🍾","玻璃瓶","Glass bottle",RECYCLE),
    ("🧴","洗发水瓶","Shampoo bottle",RECYCLE),
]
def card_grid(cards,title,sub,color,show_answer=False):
    s=ns();bg(s,CREAM)
    page_title(s,title,sub,color)
    gx,gy=0.55,1.55
    cw,ch=3.6,2.75
    gapx,gapy=0.20,0.20
    for idx,(em,cn,en,cl) in enumerate(cards):
        r=idx//2;c=idx%2
        x=gx+c*(cw+gapx);y=gy+r*(ch+gapy)
        card=dashed_rect(s,x,y,cw,ch,c=LGRAY,fill=WHITE)
        tb(s,x+0.1,y+0.35,cw-0.2,1.4,em,sz=90,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+1.95,cw-0.2,0.45,cn,sz=22,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+2.40,cw-0.2,0.3,en,sz=12,c=GRAY,a=PP_ALIGN.CENTER)
        if show_answer:
            dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+cw-0.55),Inches(y+0.12),Inches(0.4),Inches(0.4))
            dot.fill.solid();dot.fill.fore_color.rgb=cl;dot.line.color.rgb=WHITE;dot.line.width=Pt(1.5)
    return s

card_grid(sort_cards[:6],"🃏 3R 分类卡 (1/2)  Sorting Cards","剪开, 学生抽卡放进对的 R 区  ·  Cut apart — sort into the right R zone",REUSE)
card_grid(sort_cards[6:],"🃏 3R 分类卡 (2/2)  Sorting Cards","剪开, 学生抽卡放进对的 R 区  ·  Cut apart — sort into the right R zone",REUSE)

# ============================================================
# PAGES 4-5 — Detective scene cards (4 per page)
# ============================================================
detective=[
    ("🚰","水龙头一直开着","Tap left running"),
    ("🍽️","用一次就扔的餐具","Single-use plates"),
    ("💡","没有人却开着灯","Lights on, nobody there"),
    ("📄","纸只用了一面","Paper used one side"),
    ("🛍️","买东西用好多塑料袋","Too many plastic bags"),
    ("🥤","喝饮料用塑料吸管","Plastic straw + cup"),
    ("🗑️","可回收的扔进垃圾桶","Recyclables in trash"),
    ("🔌","充完电还插着","Charger left plugged in"),
]
def detective_grid(cards,part):
    s=ns();bg(s,CREAM)
    page_title(s,f"🕵️ 大侦探场景卡 ({part})  Detective Scenes",
               "找问题 → 说更好的做法  ·  Spot the problem → say a better way",REDUCE)
    gx,gy=0.55,1.55
    cw,ch=3.6,4.15
    gapx,gapy=0.20,0.20
    for idx,(em,cn,en) in enumerate(cards):
        r=idx//2;c=idx%2
        x=gx+c*(cw+gapx);y=gy+r*(ch+gapy)
        card=dashed_rect(s,x,y,cw,ch,c=LGRAY,fill=WHITE)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.15),Inches(y+0.15),Inches(0.5),Inches(0.5))
        badge.fill.solid();badge.fill.fore_color.rgb=REDUCE;badge.line.fill.background()
        tb(s,x+0.15,y+0.20,0.5,0.4,"🔍",sz=16,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+0.75,cw-0.2,1.3,em,sz=84,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+2.05,cw-0.2,0.5,cn,sz=20,b=True,c=DARK,a=PP_ALIGN.CENTER)
        tb(s,x+0.1,y+2.55,cw-0.2,0.3,en,sz=11,c=GRAY,a=PP_ALIGN.CENTER)
        # write lines
        tb(s,x+0.2,y+2.95,cw-0.4,0.3,"问题 Problem:",sz=11,b=True,c=ALERT)
        ln=s.shapes.add_connector(1,Inches(x+0.2),Inches(y+3.35),Inches(x+cw-0.2),Inches(y+3.35))
        ln.line.color.rgb=LGRAY;ln.line.width=Pt(0.75)
        tb(s,x+0.2,y+3.45,cw-0.4,0.3,"更好 Better:",sz=11,b=True,c=ECO)
        ln2=s.shapes.add_connector(1,Inches(x+0.2),Inches(y+3.85),Inches(x+cw-0.2),Inches(y+3.85))
        ln2.line.color.rgb=LGRAY;ln2.line.width=Pt(0.75)
    return s

detective_grid(detective[:4],"1/2")
detective_grid(detective[4:],"2/2")

# ============================================================
# PAGE 6 — Detective record sheet
# ============================================================
s=ns();bg(s,CREAM)
page_title(s,"📝 我是 Zero Waste 大侦探  Detective Record","名字 Name: __________________    日期 Date: __________",REDUCE)
tb(s,0.55,1.5,7.4,0.4,"看图片, 找出问题, 写下更好的做法!  Find the problem, write a better way.",sz=13,b=True,c=DARK)
# table header
hy=2.05
hdr=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.55),Inches(hy),Inches(7.4),Inches(0.5))
hdr.fill.solid();hdr.fill.fore_color.rgb=REDUCE;hdr.line.color.rgb=WHITE;hdr.line.width=Pt(1)
tb(s,0.65,hy+0.08,0.7,0.35,"#",sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(s,1.35,hy+0.08,3.0,0.35,"问题是什么? Problem",sz=13,b=True,c=WHITE)
tb(s,4.6,hy+0.08,3.3,0.35,"更好的做法? Better way + 哪个 R",sz=13,b=True,c=WHITE)
rowh=0.92
for i in range(8):
    y=hy+0.5+i*rowh
    row=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.55),Inches(y),Inches(7.4),Inches(rowh))
    row.fill.solid();row.fill.fore_color.rgb=(WHITE if i%2==0 else WARM);row.line.color.rgb=LGRAY;row.line.width=Pt(0.75)
    numc=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.72),Inches(y+0.26),Inches(0.4),Inches(0.4))
    numc.fill.solid();numc.fill.fore_color.rgb=REDUCE;numc.line.fill.background()
    tb(s,0.72,y+0.31,0.4,0.3,str(i+1),sz=13,b=True,c=WHITE,a=PP_ALIGN.CENTER)
    # divider between the two write columns
    dl=s.shapes.add_connector(1,Inches(4.5),Inches(y),Inches(4.5),Inches(y+rowh))
    dl.line.color.rgb=LGRAY;dl.line.width=Pt(0.5)

# ============================================================
# PAGE 7 — Answer key (teacher)
# ============================================================
s=ns();bg(s,CREAM)
page_title(s,"✅ 老师参考答案  Answer Key","答案不唯一 — 学生能说出理由即可  ·  Answers may vary; reasoning matters",GOLD)
# 3R sorting answers
tb(s,0.55,1.45,7.4,0.4,"🃏 3R 分类卡 (蓝=Reduce · 橙=Reuse · 绿=Recycle)",sz=15,b=True,c=ECO)
ans=[("💧 一直流的水龙头","Reduce · 用完就关",REDUCE),
     ("📄 只写一面的纸","Reduce · 两面都用",REDUCE),
     ("🛍️ 塑料袋","Reduce · 自带布袋",REDUCE),
     ("🥤 塑料吸管","Reduce · 自带水杯",REDUCE),
     ("🍶 空塑料瓶","Reuse/Recycle · 花盆或回收",REUSE),
     ("👕 小了的衣服","Reuse · 传给弟妹",REUSE),
     ("📦 快递纸箱","Reuse/Recycle · 做手工或回收",REUSE),
     ("🫙 玻璃罐子","Reuse · 当笔筒",REUSE),
     ("🥫 喝完的易拉罐","Recycle · 金属回收",RECYCLE),
     ("📰 旧报纸","Recycle · 纸类回收",RECYCLE),
     ("🍾 玻璃瓶","Recycle · 玻璃回收",RECYCLE),
     ("🧴 洗发水瓶","Recycle · 塑料回收",RECYCLE)]
for i,(item,a,cl) in enumerate(ans):
    r=i%6;c=i//6
    x=0.55+c*3.85;y=1.9+r*0.52
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y+0.05),Inches(0.22),Inches(0.22))
    dot.fill.solid();dot.fill.fore_color.rgb=cl;dot.line.fill.background()
    tb(s,x+0.32,y,3.45,0.45,item,sz=11,b=True,c=DARK)
    tb(s,x+0.32,y+0.26,3.45,0.25,a,sz=9,c=GRAY)
# Detective answers
tb(s,0.55,5.2,7.4,0.4,"🕵️ 大侦探场景 — 问题 → 更好的做法 (R)",sz=15,b=True,c=REDUCE)
dans=[("🚰 水龙头一直开","关掉水龙头 (Reduce)"),
      ("🍽️ 一次性餐具","用可洗的餐具 (Reduce/Reuse)"),
      ("💡 没人还开灯","离开就关灯 (Reduce)"),
      ("📄 纸只用一面","两面用 / 回收 (Reduce/Recycle)"),
      ("🛍️ 好多塑料袋","自带布袋 (Reduce/Reuse)"),
      ("🥤 塑料吸管","自带水杯 (Reduce)"),
      ("🗑️ 可回收扔垃圾桶","放进回收箱 (Recycle)"),
      ("🔌 充完还插着","拔掉充电器 (Reduce)")]
for i,(q,a) in enumerate(dans):
    r=i%4;c=i//4
    x=0.55+c*3.85;y=5.7+r*0.62
    tb(s,x,y,3.7,0.35,q,sz=11,b=True,c=DARK)
    tb(s,x,y+0.26,3.7,0.3,"→ "+a,sz=10,c=ECO)
# footer tip
ft=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.55),Inches(8.5),Inches(7.4),Inches(0.85))
ft.fill.solid();ft.fill.fore_color.rgb=WARM;ft.line.color.rgb=ECO;ft.line.width=Pt(1.5)
tb(s,0.7,8.6,7.1,0.35,"💡 教学优先级: 能 Reduce 就 Reduce → 不能才 Reuse → 最后才 Recycle。",sz=12,b=True,c=ECO)
tb(s,0.7,8.95,7.1,0.3,"Priority: Reduce first, then Reuse, Recycle last.",sz=11,c=GRAY)

# === Save ===
out=os.path.join(BASE,"zerowaste_3R_cards.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} pages)")
