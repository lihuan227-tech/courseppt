#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_reskinned_deck.py  —  Zero Waste · Day 1
Takes the CONTENT + structure of day1_trash_clean_deck.pptx (build_clean_deck.py)
and re-skins every slide in the DESIGN LANGUAGE of day1_trash.pptx (_helpers.py):

  • KaiTi font, airy spaced-CJK header titles
  • sage-cream content backgrounds
  • signature full-width COLORED HEADER BAR with the page number tucked inside it
  • navy (INK) + gold (STAR) cover, star-studded section dividers
  • INK + STAR sentence-frame / key-idea / exit panels

Strategy: import build_clean_deck as the content engine, then monkey-patch its
chrome (font, base, header, cover, divider, group photo) and its framed panels
(s_frame / s_statement / s_keyidea / s_share / s_practice) to A's look. All other
slide builders inherit the new chrome + font automatically and keep B's content.

Output: day1_trash_reskinned.pptx   (NEW file)
Run:    python3 build_reskinned_deck.py
Deps:   python-pptx
"""
import os
import build_clean_deck as D
from build_clean_deck import H, mix, T, Box, IN, rect, text, card, pill
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "day1_trash_reskinned.pptx")

# ---------------------------------------------------------------- A palette
INK   = H('0F1A3A')   # deep navy — cover / dividers accents / frame panels
STAR  = H('F5C242')   # gold accent on dark
WARM  = H('FFF3E0')
CREAMA= H('F6F8EC')   # sage-tinted cream page bg
WHITE = D.WHITE
GRAY  = D.GRAY
LGRAY = H('BBBBBB')
DARK  = H('2C2C2C')
GREEN = D.GREEN
MOSS  = D.MOSS
SW, SH = D.SW, D.SH
FONT = 'KaiTi'

# ---- global font + cream swap (read at call-time by every D.text/base) ----
D.CJK = FONT
D.DISP = FONT
D.CREAM = CREAMA

# ---------------------------------------------------------------- helpers
def _cjk(ch):
    return '一' <= ch <= '鿿'

def sp(t):
    """Airy spaced-CJK like the day1_trash design (thin space between hanzi)."""
    out = []
    for i, ch in enumerate(t):
        out.append(ch)
        if i + 1 < len(t) and _cjk(ch) and _cjk(t[i + 1]):
            out.append(' ')  # thin space
    return ''.join(out)

# ============================================================ CHROME
def base(slide, accent=GREEN, section_tab="", bg=None):
    rect(slide, Box(0, 0, SW, SH), CREAMA, radius=0)
    D.PAGE[0] += 1
    text(slide, Box(0.55, SH - 0.42, 7, 0.3),
         [[T("谷雨中文 GR EDU · 零废弃 Zero Waste · Day 1", 10, GRAY, False, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE)

def section_tab_pill(slide, label, accent):
    return  # A design keeps the page number in the header bar instead

def header(slide, kicker, title_zh, title_en="", accent=GREEN):
    """A-style full-width colored header bar + page number inside it."""
    bar = Box(0.45, 0.26, 12.43, 0.86)
    rect(slide, bar, accent, radius=0.16)
    runs = [T(sp(title_zh), 23, WHITE, True, False, FONT)]
    if title_en:
        runs.append(T("   " + title_en, 13, mix(tuple(WHITE), tuple(accent), .85), False, True, FONT))
    text(slide, Box(bar.x + 0.28, bar.y, bar.w - 1.2, bar.h), [runs],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Box(SW - 1.75, bar.y, 1.0, bar.h),
         [[T(f"{D.PAGE[0]}", 13, WHITE, True, False, FONT)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    top = bar.y + bar.h + 0.12
    if kicker:
        text(slide, Box(0.75, top, 11.5, 0.32),
             [[T(kicker.upper(), 12, mix(tuple(accent), (255, 255, 255), .55), True, False, FONT)]])
        top += 0.40
    return top

def frame_bar(slide, box, frame_cn, frame_en, label="💬 句型"):
    """A signature sentence-frame banner: INK fill + STAR border + gold text."""
    rect(slide, box, INK, line=STAR, line_w=2, radius=0.18)
    paras = [[T(f"{label}: ", 13, STAR, True, False, FONT),
              T(frame_cn, 15, WHITE, True, False, FONT)]]
    if frame_en:
        paras.append([T(frame_en, 10.5, WARM, False, True, FONT)])
    text(slide, box.pad(0.12), paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)

# ---------------------------------------------------------------- COVER (A)
def s_cover(prs):
    s = D.new_slide(prs)
    rect(s, Box(0, 0, SW, SH), INK, radius=0)
    # faint accent circles
    for cx, cy, d, col in [(10.8, -2.8, 8.4, mix(tuple(GREEN), tuple(INK), .55)),
                           (-2.6, 5.4, 6.6, mix(tuple(GREEN), tuple(INK), .42))]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(cx), IN(cy), IN(d), IN(d))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background(); c.shadow.inherit = False
    text(s, Box(0.6, 0.55, 12.1, 0.5),
         [[T(sp("♻️ 零废弃与可持续发展"), 18, STAR, True, False, FONT),
           T("   Zero Waste & Sustainability", 14, STAR, False, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0.6, 1.55, 12.1, 1.0),
         [[T(sp("Day 1 · 垃圾去哪儿了?"), 44, WHITE, True, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0.6, 2.7, 12.1, 0.45),
         [[T("Where Does Our Trash Go?", 18, LGRAY, False, True, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0.6, 3.35, 12.1, 0.9), [[T("🗑️   🍌   📰   🥤", 52, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    q = Box(2.4, 4.5, 8.5, 1.9)
    rect(s, q, GREEN, line=STAR, line_w=3, radius=0.08)
    text(s, Box(q.x, q.y + 0.18, q.w, 0.4),
         [[T(sp("🤔 今天的探究问题  Today's Inquiry"), 14, STAR, True, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(q.x, q.y + 0.66, q.w, 0.5),
         [[T(sp("你扔的垃圾,要放进哪个桶?"), 20, WHITE, True, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(q.x, q.y + 1.18, q.w, 0.5),
         [[T("⚠ 一起打败「垃圾怪」!  Which bin does your trash go in? Defeat the Trash Monster!",
             11.5, WARM, False, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0.6, SH - 0.5, 12.1, 0.3),
         [[T("谷雨中文 GR EDU · 暑期中文 · 社会科学 / 科学", 11, mix(tuple(WHITE), tuple(INK), .7), False, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- DIVIDER (A)
def s_divider(prs, sect, zh, en, meta=""):
    s = D.new_slide(prs); a = D.SECT[sect]['accent']
    rect(s, Box(0, 0, SW, SH), a, radius=0)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(9.6), IN(3.4), IN(7.2), IN(7.2))
    c.fill.solid(); c.fill.fore_color.rgb = mix(tuple(WHITE), tuple(a), .14)
    c.line.fill.background(); c.shadow.inherit = False
    text(s, Box(0.9, 1.55, 11.5, 0.5),
         [[T(D.SECT[sect]['no'], 16, STAR, True, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0.9, 2.15, 11.5, 1.4), [[T(sp(zh), 40, WHITE, True, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0.92, 3.7, 11.5, 0.6), [[T(en, 18, mix(tuple(WHITE), tuple(a), .85), False, True, FONT)]])
    if meta:
        pill(s, Box(0.92, 4.5, min(0.14 * len(meta) + 1.2, 6.5), 0.6),
             mix(tuple(WHITE), tuple(a), .18), meta, 15, WHITE, bold=True, fnt=FONT)
    # gold stars
    for x, y, sz in [(1.0, 6.3, 0.42), (1.9, 6.05, 0.34), (10.9, 6.05, 0.34), (11.7, 6.3, 0.42)]:
        st = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, IN(x), IN(y), IN(sz), IN(sz))
        st.fill.solid(); st.fill.fore_color.rgb = STAR; st.line.fill.background(); st.shadow.inherit = False

# ---------------------------------------------------------------- GROUP PHOTO (A)
def s_groupphoto(prs):
    s = D.new_slide(prs)
    rect(s, Box(0, 0, SW, SH), GREEN, radius=0)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(9.0), IN(-2.0), IN(7), IN(7))
    c.fill.solid(); c.fill.fore_color.rgb = mix(tuple(MOSS), tuple(GREEN), .5)
    c.line.fill.background(); c.shadow.inherit = False
    text(s, Box(0, 1.8, SW, 1.2), [[T("📸", 60, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0, 3.0, SW, 1.0), [[T(sp("全班合影 — 拿着作品!"), 40, WHITE, True, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0, 4.15, SW, 0.5), [[T("Group photo — hold up your work! 🎉", 17, STAR, False, True, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Box(0, 4.95, SW, 0.8), [[T("🎡   🔑   ♻️   🌍   🌱", 36, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for x, y, sz in [(1.2, 6.3, 0.4), (2.0, 6.05, 0.32), (10.9, 6.05, 0.32), (11.6, 6.3, 0.4)]:
        st = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, IN(x), IN(y), IN(sz), IN(sz))
        st.fill.solid(); st.fill.fore_color.rgb = STAR; st.line.fill.background(); st.shadow.inherit = False

# ============================================================ FRAMED PANELS → INK + STAR
def s_frame(prs, sect, tag, zh_html, en, note=""):
    s = D.new_slide(prs); a = D.SECT[sect]['accent']
    base(s, a); top = header(s, "句型 · sentence frame", "会用句型说一说", "Say it with the sentence frame", a)
    area = Box(0.9, top + 0.2, SW - 1.8, SH - top - 1.0)
    b = Box(area.x, area.y + 0.2, area.w, 2.4)
    rect(s, b, INK, line=STAR, line_w=2.5, radius=0.06, shadow=True)
    pill(s, Box(b.x + 0.4, b.y + 0.4, 2.0, 0.6), STAR, tag, 14, INK, bold=True, fnt=FONT)
    text(s, Box(b.x + 0.4, b.y + 1.15, b.w - 0.8, 1.0), [[T(sp(zh_html), 30, WHITE, True, False, FONT)]])
    text(s, Box(b.x + 0.4, b.y + 1.85, b.w - 0.8, 0.5), [[T(en, 15, STAR, False, True, FONT)]])
    if note:
        pill(s, Box(area.x, b.y + b.h + 0.35, min(0.16 * len(note) + 1.2, 8), 0.6),
             D.SOFT, note, 13, GREEN, bold=True, fnt=FONT)

def s_statement(prs, sect, kicker, zh, en, frame_zh, frame_en, closing=None):
    s = D.new_slide(prs); a = D.SECT[sect]['accent']
    base(s, a); top = header(s, kicker, zh, en, a)
    area = Box(0.9, top + 0.3, SW - 1.8, SH - top - 1.0)
    fb = Box(area.x + 0.5, area.y + 0.1, area.w - 1.0, 1.95)
    rect(s, fb, INK, line=STAR, line_w=3, radius=0.06, shadow=True)
    text(s, fb, [[T(sp(frame_zh), 30, WHITE, True, False, FONT)], [T(frame_en, 14, STAR, False, True, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.4)
    D.bin_chips_row(s, Box(area.x + 1.4, fb.y + fb.h + 0.4, area.w - 2.8, 0.6))
    if closing:
        pill(s, Box(SW / 2 - 3.4, fb.y + fb.h + 1.2, 6.8, 0.62), GREEN, closing, 15, STAR, bold=True, fnt=FONT)

def s_keyidea(prs, sect, kicker, zh, en, qs, idea_zh, idea_en):
    s = D.new_slide(prs); a = D.SECT[sect]['accent']
    base(s, a); top = header(s, kicker, zh, en, a)
    area = Box(0.9, top + 0.15, SW - 1.8, SH - top - 0.8)
    lw = area.w * 0.5
    n = len(qs); gap = 0.24; qh = (area.h - gap * (n - 1)) / n
    for i, (q, qe) in enumerate(qs):
        b = Box(area.x, area.y + i * (qh + gap), lw, qh); card(s, b)
        rect(s, Box(b.x + 0.18, b.y + qh / 2 - 0.26, 0.52, 0.52), a, radius=0.3)
        text(s, Box(b.x + 0.18, b.y + qh / 2 - 0.26, 0.52, 0.52), [[T(str(i + 1), 18, WHITE, True, False, FONT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        D.bilingual(s, Box(b.x + 0.9, b.y, b.w - 1.1, qh), q, qe, zh_sz=17, en_sz=11.5)
    pb = Box(area.x + lw + 0.35, area.y, area.w - lw - 0.35, area.h)
    rect(s, pb, INK, line=STAR, line_w=2.5, radius=0.06, shadow=True)
    text(s, Box(pb.x + 0.35, pb.y + 0.3, pb.w - 0.7, 0.5), [[T("💡 小结 Key idea", 15, STAR, True, False, FONT)]])
    text(s, Box(pb.x + 0.35, pb.y + 0.95, pb.w - 0.7, pb.h - 1.4),
         [[T(sp(idea_zh), 22, WHITE, True, False, FONT)], [T(idea_en, 13, STAR, False, True, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, spacing=1.3)

def s_share(prs, kicker, zh, en, tag, fz, fe, emoji, sub, accent=D.ORANGE):
    s = D.new_slide(prs); a = accent
    base(s, a); top = header(s, kicker, zh, en, a)
    area = Box(0.9, top + 0.3, SW - 1.8, SH - top - 1.0)
    fb = Box(area.x, area.y + 0.1, area.w, 1.7)
    rect(s, fb, INK, line=STAR, line_w=2.5, radius=0.06, shadow=True)
    pill(s, Box(fb.x + 0.4, fb.y + fb.h / 2 - 0.28, 2.0, 0.56), STAR,
         tag, 14, INK, bold=True, fnt=FONT)
    text(s, Box(fb.x + 2.7, fb.y, fb.w - 3.0, fb.h),
         [[T(sp(fz), 25, WHITE, True, False, FONT)], [T(fe, 13, STAR, False, True, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, spacing=1.3)
    cb = Box(area.x + area.w / 2 - 3.2, fb.y + fb.h + 0.35, 6.4, area.h - fb.h - 0.45)
    card(s, cb)
    text(s, cb, [[T(emoji, 46, DARK)], [T(sp(sub), 18, DARK, True, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)

def s_practice(prs, sect, i, total, item, en, emoji, binidx, why, frame_zh, frame_en):
    # QUESTION mode — no answer shown (the answer key comes on a later slide)
    s = D.new_slide(prs); a = D.SECT[sect]['accent']
    base(s, a); top = header(s, f"放进哪个桶? {i}/{total} · which bin?",
                             f"{item},放进哪个桶?", f"Which bin for a {en}?", a)
    area = Box(0.9, top + 0.1, SW - 1.8, SH - top - 0.85)
    lw = 3.7; lb = Box(area.x, area.y, lw, area.h)
    card(s, lb, shadow=True)
    ph = Box(lb.x + 0.35, lb.y + 0.35, lb.w - 0.7, lb.h * 0.5)
    rect(s, ph, H('EAF0E2'), line=mix(tuple(a), (255, 255, 255), .4), line_w=1.5, radius=0.06, dash=True)
    text(s, ph, [[T("📷", 34, mix(tuple(a), (255, 255, 255), .3))],
                 [T("实物照片 · insert photo", 11, GRAY, False, True, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
    text(s, Box(lb.x, lb.y + lb.h * 0.5 + 0.5, lb.w, lb.h * 0.5 - 0.6),
         [[T(sp(item), 26, DARK, True, False, FONT)], [T(en, 12, GRAY, False, True, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, spacing=1.15)
    rb = Box(lb.x + lw + 0.4, area.y, area.w - lw - 0.4, area.h)
    grid_h = rb.h * 0.52
    n = 4; gap = 0.2; cw = (rb.w - gap * (n - 1)) / n
    for j, b2 in enumerate(D.BINS):
        mb = Box(rb.x + j * (cw + gap), rb.y, cw, grid_h)
        rect(s, mb, WHITE, line=b2['c'], line_w=2, radius=0.12, shadow=True)
        rect(s, Box(mb.x, mb.y, mb.w, 0.46), b2['c'], radius=0.12)
        rect(s, Box(mb.x, mb.y + 0.25, mb.w, 0.21), b2['c'], radius=0)
        text(s, Box(mb.x, mb.y + 0.66, mb.w, grid_h - 0.7), [[T(sp(b2['zh']), 13, b2['c'], True, False, FONT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pb = Box(rb.x, rb.y + grid_h + 0.28, rb.w, 0.58)
    rect(s, pb, D.SOFT, radius=0.16)
    text(s, pb, [[T("🤔 想一想:它放哪个桶? 先猜一猜 — 等下看答案!", 14, GREEN, True, False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    frame_bar(s, Box(rb.x, rb.y + grid_h + 1.0, rb.w, 0.62),
              "这是 ____。我猜它放进 ____ 桶。", "This is ___. I guess it goes in the ___ bin.", label="我来猜")

# ============================================================ wire up overrides
D.base = base
D.section_tab_pill = section_tab_pill
D.header = header
D.s_cover = s_cover
D.s_divider = s_divider
D.s_groupphoto = s_groupphoto
D.s_frame = s_frame
D.s_statement = s_statement
D.s_keyidea = s_keyidea
D.s_share = s_share
D.s_practice = s_practice
D.OUT = OUT

if __name__ == '__main__':
    D.build()
    print("RESKINNED →", OUT)
