#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
patch_day3_water.py — add 2 slides to PPT/day3_water_plastic.pptx IN PLACE,
preserving all embedded photos (the base build script has no image pipeline,
so this deck is edited directly rather than regenerated).

Adds, after the existing Bamboozle slide (index 29):
  • a separate "Review: Bamboozle" page (blank link placeholder)
  • a "先完成练习册 Complete Your Booklet" slide (Day 1–3)
This deck has no projects section, so the booklet slide sits before the wrap-up.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "PPT", "day3_water_plastic.pptx")

# --- Ocean palette (match create_day2_water.py) ---
DEEP  = RGBColor(0x0B,0x55,0x63); OCEAN = RGBColor(0x15,0x65,0xA0)
SEAGRN= RGBColor(0x2E,0x8B,0x7A); CREAM = RGBColor(0xFB,0xF7,0xEC)
CORAL = RGBColor(0xE0,0x63,0x3F); SUNYEL= RGBColor(0xF5,0xC2,0x42)
WHITE = RGBColor(0xFF,0xFF,0xFF); DARK  = RGBColor(0x2C,0x2C,0x2C)
GRAY  = RGBColor(0x88,0x88,0x88); WARM  = RGBColor(0xFF,0xF3,0xE0)
OK    = RGBColor(0x38,0x8E,0x3C)
FONT='KaiTi'

prs = Presentation(DECK)
W,H = prs.slide_width, prs.slide_height

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])
def tb(s,l,t,w,h,txt,sz=18,b=False,c=DARK,a=None):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    if a: p.alignment=a
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=FONT
    return tf
def bg(s,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    sp=sh._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
def box(s,l,t,w,h,fill=WHITE,line=None,lw=2.0,rad=True):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False; return sh
def hb(s,txt,c=DEEP,t=0.18):
    box(s,0.3,t,9.4,0.55,fill=c); tb(s,0.45,t+0.07,9.1,0.45,txt,sz=16,b=True,c=WHITE)

# ---- Slide A: Review Bamboozle (blank link) ----
sA=ns(); bg(sA,CREAM); hb(sA,"🎮 复习 · Bamboozle 大复习  Review",DEEP)
tb(sA,0.4,0.85,9.2,0.3,"复习游戏 — 分组抢答!",13,True,DARK,PP_ALIGN.CENTER)
box(sA,0.4,1.2,4.4,3.6,fill=DEEP)
tb(sA,0.4,1.5,4.4,1.0,"🎮",72,a=PP_ALIGN.CENTER)
box(sA,1.05,2.8,3.1,0.6,fill=CORAL)
tb(sA,1.05,2.9,3.1,0.42,"▶️ 开始复习",15,True,WHITE,PP_ALIGN.CENTER)
box(sA,0.7,3.95,3.8,0.55,fill=WHITE,line=SUNYEL,lw=1.5)
tb(sA,0.85,4.05,3.55,0.38,"🔗 链接: ____________ (老师粘贴)",11,True,GRAY,PP_ALIGN.LEFT)
box(sA,5.1,1.2,4.5,3.6,fill=WHITE,line=OCEAN,lw=2.5)
tb(sA,5.3,1.35,4.1,0.4,"📋 怎么玩 How to Play",15,True,OCEAN)
for i,t in enumerate(["1. 老师点开上面的链接","2. 全班分 2-3 组","3. 轮流抢答, 答对加分","4. 复习今天: 水 + 塑料 + 海洋小卫士"]):
    tb(sA,5.35,1.95+i*0.6,4.1,0.5,t,13,True,DARK)

# ---- Slide B: Complete Booklet (Day 1–3) ----
sB=ns(); bg(sB,CREAM); hb(sB,"📓 先完成练习册!  Booklet First",DEEP)
tb(sB,0.4,0.85,9.2,0.3,"收尾之前, 先完成今天的练习册!",13,True,DARK,PP_ALIGN.CENTER)
books=[("Day 1","垃圾分类",SEAGRN),("Day 2","可再生能源",SUNYEL),("Day 3","水与塑料",OCEAN)]
bw=2.7; bgap=0.45; bstart=(10-3*bw-2*bgap)/2
for i,(lab,zh,cl) in enumerate(books):
    x=bstart+i*(bw+bgap)
    box(sB,x,1.3,bw,3.1,fill=WHITE,line=cl,lw=2.5)
    tb(sB,x,1.6,bw,0.9,"📓",50,a=PP_ALIGN.CENTER)
    tb(sB,x,2.6,bw,0.5,lab,20,True,cl,PP_ALIGN.CENTER)
    tb(sB,x,3.12,bw,0.35,zh,13,True,DARK,PP_ALIGN.CENTER)
    box(sB,x+bw/2-1.0,3.65,2.0,0.5,fill=WARM,line=cl,lw=1.2)
    tb(sB,x+bw/2-1.0,3.73,2.0,0.36,"✅ 完成 → 打勾",12,True,OK,PP_ALIGN.CENTER)
box(sB,0.4,4.7,9.2,0.55,fill=WARM,line=SUNYEL,lw=1.5)
tb(sB,0.55,4.8,9.0,0.36,"👩‍🏫 先看 → 一起读 → 自己写 → 同桌检查",12,True,DARK,PP_ALIGN.CENTER)

# ---- reorder: move the 2 new slides to just after the Bamboozle (index 29) ----
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
review_el, booklet_el = ids[-2], ids[-1]
sldIdLst.remove(review_el); sldIdLst.remove(booklet_el)
anchor = ids[29]            # the existing Bamboozle slide
anchor.addnext(review_el)
review_el.addnext(booklet_el)

prs.save(DECK)
print("PATCHED", DECK, "->", len(prs.slides._sldIdLst), "slides")
