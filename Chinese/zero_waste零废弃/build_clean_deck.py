#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_clean_deck.py  —  Zero Waste · Day 1  ·  CLEAN content-first rebuild

NOT a port of day1_trash_version1.html. This authors the same teaching content
(4 sessions: sort the 4 bins → words & writing → recycled crafts → wrap-up) into
a fresh, disciplined "clean kid-friendly" design system: one consistent header,
generous whitespace, a small reusable layout vocabulary, real visual hierarchy.

Output: day1_trash_clean_deck.pptx   (NEW file)
Run:    python3 build_clean_deck.py
Deps:   python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE=os.path.dirname(os.path.abspath(__file__))
OUT=os.path.join(HERE,"day1_trash_clean_deck.pptx")

# ---------------------------------------------------------------- palette
def H(h):
    h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
def mix(a,b,p):  # p of a
    return RGBColor(*[round(x*p+y*(1-p)) for x,y in zip(tuple(a),tuple(b))])
WHITE=H('FFFFFF')
INK=H('27313A')          # softer than pure navy for body text
GREEN=H('2C5F2D')        # forest — primary
MOSS=H('6BA03D')
TEAL=H('006970')
ORANGE=H('E65C00')
CREAM=H('F7F8F0')        # flat page bg
PANEL=H('FFFFFF')
LINE=H('E4E2D4')         # subtle card border
GRAY=H('8A8F86')         # captions / english
SOFT=H('F0F2E6')         # soft fill chips
STAR=H('F5B921')
# four bins (official standard)
B_RECYCLE=H('1B6FBA'); B_FOOD=H('5C9A35'); B_HARM=H('C8253E'); B_OTHER=H('7C8088')

SW,SH=13.333,7.5
def IN(v): return Emu(int(round(v*914400)))

CJK="Noto Sans SC"   # editable, CJK-safe; PowerPoint will substitute gracefully
DISP="Noto Sans SC"

class Box:
    __slots__=('x','y','w','h')
    def __init__(s,x,y,w,h): s.x,s.y,s.w,s.h=float(x),float(y),float(w),float(h)
    def pad(s,p): return Box(s.x+p,s.y+p,s.w-2*p,s.h-2*p)

# ---------------------------------------------------------------- primitives
def _ea(run,font):
    rPr=run._r.get_or_add_rPr()
    for tag in ('a:ea','a:cs'):
        e=rPr.find(qn(tag))
        if e is None:
            e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set('typeface',font)

def text(slide,box,paras,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
         wrap=True,spacing=1.06,space_after=2.0):
    """paras: list of paragraphs; paragraph = list of run dicts/tuples
       run tuple: (txt,size,color,bold,italic,font)"""
    tb=slide.shapes.add_textbox(IN(box.x),IN(box.y),IN(box.w),IN(box.h))
    tf=tb.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=IN(0.04); tf.margin_right=IN(0.04)
    tf.margin_top=IN(0.02); tf.margin_bottom=IN(0.02)
    first=True
    for para in paras:
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.alignment=align; p.line_spacing=spacing
        p.space_after=Pt(space_after); p.space_before=Pt(0)
        for (t,sz,col,bold,ital,fnt) in para:
            r=p.add_run(); r.text=t
            r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=ital
            r.font.color.rgb=col; r.font.name=fnt or CJK; _ea(r,fnt or CJK)
    return tb

def T(t,sz,col,bold=False,ital=False,fnt=None): return (t,sz,col,bold,ital,fnt)

def rect(slide,box,fill,line=None,line_w=1.0,radius=0.10,shadow=False,dash=False):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if radius>0 else MSO_SHAPE.RECTANGLE
    sp=slide.shapes.add_shape(shp,IN(box.x),IN(box.y),IN(box.w),IN(box.h))
    if radius>0:
        try: sp.adjustments[0]=radius
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=line; sp.line.width=Pt(line_w)
        if dash:
            ln=sp.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    sp.shadow.inherit=False
    if shadow:
        _soft_shadow(sp)
    sp.text_frame.word_wrap=True
    return sp

def _soft_shadow(sp):
    spPr=sp._element.spPr
    el=spPr.makeelement(qn('a:effectLst'),{})
    sh=spPr.makeelement(qn('a:outerShdw'),
        {'blurRad':'90000','dist':'38100','dir':'5400000','rotWithShape':'0'})
    clr=spPr.makeelement(qn('a:srgbClr'),{'val':'2C2C2C'})
    a=spPr.makeelement(qn('a:alpha'),{'val':'16000'})
    clr.append(a); sh.append(clr); el.append(sh); spPr.append(el)

def card(slide,box,fill=PANEL,line=LINE,radius=0.055,shadow=True,line_w=1.0):
    return rect(slide,box,fill,line=line,line_w=line_w,radius=radius,shadow=shadow)

def pill(slide,box,fill,txt,size,tcol,bold=True,fnt=None,line=None,line_w=1.0):
    rect(slide,box,fill,line=line,line_w=line_w,radius=0.5)
    text(slide,box,[[T(txt,size,tcol,bold,False,fnt)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def swatch(slide,cx,cy,size,color):
    rect(slide,Box(cx-size/2,cy-size/2,size,size),color,radius=0.22)

# ---------------------------------------------------------------- page chrome
PAGE=[0]
def base(slide,accent=GREEN,section_tab="",bg=CREAM):
    rect(slide,Box(0,0,SW,SH),bg,radius=0)
    # left accent spine
    rect(slide,Box(0,0,0.16,SH),accent,radius=0)
    # footer
    PAGE[0]+=1
    text(slide,Box(0.7,SH-0.5,6,0.3),[[T("谷雨中文 GR EDU · 零废弃 Zero Waste · Day 1",10.5,GRAY,False,False)]],anchor=MSO_ANCHOR.MIDDLE)
    text(slide,Box(SW-1.5,SH-0.5,0.8,0.3),[[T(f"{PAGE[0]:02d}",11,GRAY,True)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    if section_tab:
        bw=min(0.18*len(section_tab)+0.7,3.4)
        b=Box(SW-bw-0.6,0.5,bw,0.42)
        pill(slide,b,mix(tuple(accent),(255,255,255),.14),section_tab,11.5,accent if accent!=GREEN else GREEN,bold=True)
        # recolor pill text white on dark
        # (simpler: redo as filled accent)
    return

def header(slide,kicker,title_zh,title_en="",accent=GREEN):
    x=0.72; y=0.62; w=SW-1.5
    if kicker:
        text(slide,Box(x,y,w,0.32),[[T(kicker.upper(),12.5,accent,True)]])
        y+=0.40
    text(slide,Box(x,y,w,0.72),[[T(title_zh,29,GREEN,True,False,DISP)]])
    y+=0.70
    # accent underline
    rect(slide,Box(x+0.02,y+0.02,0.95,0.045),accent,radius=0.5)
    if title_en:
        text(slide,Box(x+1.1,y-0.06,w-1.1,0.3),[[T(title_en,12.5,GRAY,False,True)]],anchor=MSO_ANCHOR.MIDDLE)
    return y+0.30   # content-top y

def section_tab_pill(slide,label,accent):
    bw=min(0.16*len(label)+0.8,3.6)
    b=Box(SW-bw-0.6,0.6,bw,0.44)
    rect(slide,b,accent,radius=0.5)
    text(slide,b,[[T(label,11.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# bilingual helper: zh primary + en caption block inside a region
def bilingual(slide,box,zh,en,zh_sz=17,en_sz=11.5,zh_col=INK,bold=True,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE):
    paras=[[T(zh,zh_sz,zh_col,bold)]]
    if en: paras.append([T(en,en_sz,GRAY,False,True)])
    text(slide,box,paras,align=align,anchor=anchor,spacing=1.1)

print("helpers loaded")

# ================================================================ section meta
SECT={
 1:{'accent':GREEN,'tab':'第一节 · 分类','no':'SESSION 1'},
 2:{'accent':B_RECYCLE,'tab':'第二节 · 词语','no':'SESSION 2'},
 3:{'accent':ORANGE,'tab':'第三节 · 手作','no':'SESSION 3'},
 4:{'accent':MOSS,'tab':'延伸 · 总结','no':'WRAP-UP'},
}
BINS=[
 {'c':B_RECYCLE,'ic':'♻️','zh':'可回收物','py':'kě huí shōu wù','en':'Recyclable','dest':'回收 → 做成新东西','desten':'Recycle into new things',
  'like':'「我喜欢可以再用的东西!」','liken':'I like things we can use again!',
  'ex':[('🫙','玻璃罐','glass jar'),('📦','快递箱','box'),('📚','旧书本','old book'),('👕','旧衣服','old clothes'),('🧴','塑料瓶','bottle')]},
 {'c':B_FOOD,'ic':'🍎','zh':'厨余垃圾','py':'chú yú lā jī','en':'Food waste','dest':'堆肥 → 变成肥料','desten':'Compost into fertilizer',
  'like':'「我喜欢会腐烂的东西!」','liken':'I like things that rot!',
  'ex':[('🍚','剩饭','rice'),('🥬','剩菜','leftovers'),('🦴','骨头','bone'),('🍂','落叶','leaf'),('🥀','残花','flower')]},
 {'c':B_HARM,'ic':'⚠️','zh':'有害垃圾','py':'yǒu hài lā jī','en':'Hazardous','dest':'特殊处理 → 安全收走','desten':'Special, safe handling','harm':True,
  'like':'「我对人和环境有害!」','liken':'I can hurt people & nature!',
  'ex':[('💊','过期药品','medicine'),('🌡️','体温计','thermometer'),('🔋','废旧电池','battery'),('💡','旧灯泡','bulb')]},
 {'c':B_OTHER,'ic':'🗑️','zh':'其他垃圾','py':'qí tā lā jī','en':'Other waste','dest':'填埋 / 焚烧','desten':'Landfill / incinerate',
  'like':'「剩下的垃圾都给我!」','liken':'All the rest comes to me!',
  'ex':[('🌫️','尘土','dust'),('🧱','砖块','brick'),('🧻','脏纸巾','tissue'),('🛍️','包装袋','plastic wrap')]},
]

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ---------------------------------------------------------------- COVER
def s_cover(prs):
    s=new_slide(prs)
    rect(s,Box(0,0,SW,SH),GREEN,radius=0)
    rect(s,Box(0,0,SW,SH),GREEN,radius=0)
    # big soft circle accent
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,IN(8.6),IN(-2.2),IN(7.5),IN(7.5))
    c.fill.solid(); c.fill.fore_color.rgb=mix(tuple(MOSS),tuple(GREEN),.55); c.line.fill.background(); c.shadow.inherit=False
    c2=s.shapes.add_shape(MSO_SHAPE.OVAL,IN(-2.4),IN(4.0),IN(6.2),IN(6.2))
    c2.fill.solid(); c2.fill.fore_color.rgb=mix(tuple(MOSS),tuple(GREEN),.4); c2.line.fill.background(); c2.shadow.inherit=False
    pill(s,Box(0.95,1.15,4.7,0.6),mix(tuple(MOSS),tuple(GREEN),.5),"ZERO WASTE · 零废弃 · DAY 1",15,WHITE,bold=True)
    text(s,Box(0.9,2.15,9.5,2.0),[[T("垃圾去哪儿了?",60,WHITE,True,False,DISP)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(0.95,4.2,8.6,1.0),
         [[T("你扔的垃圾，要放进哪个桶?",24,WHITE,True)],
          [T("Which bin does your trash go in?",15,mix(tuple(WHITE),tuple(MOSS),.7),False,True)]],spacing=1.3)
    pill(s,Box(0.95,5.75,6.2,0.62),STAR,"⚠ 一起打败「垃圾怪」!  Defeat the Trash Monster!",14,GREEN,bold=True)
    text(s,Box(0.9,SH-0.55,8,0.3),[[T("谷雨中文 GR EDU · 暑期中文 · 社会科学 / 科学",11,mix(tuple(WHITE),tuple(MOSS),.75),False,False)]])

# ---------------------------------------------------------------- SECTION DIVIDER
def s_divider(prs,sect,zh,en,meta=""):
    s=new_slide(prs); a=SECT[sect]['accent']
    rect(s,Box(0,0,SW,SH),a,radius=0)
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,IN(9.2),IN(3.6),IN(6.6),IN(6.6))
    c.fill.solid(); c.fill.fore_color.rgb=mix(tuple(WHITE),tuple(a),.16); c.line.fill.background(); c.shadow.inherit=False
    text(s,Box(0.95,2.0,3,1.0),[[T(SECT[sect]['no'],16,mix(tuple(WHITE),tuple(a),.78),True)]])
    text(s,Box(0.92,2.55,10.5,1.4),[[T(zh,46,WHITE,True,False,DISP)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(0.95,4.05,10.5,0.6),[[T(en,18,mix(tuple(WHITE),tuple(a),.82),False,True)]])
    if meta:
        pill(s,Box(0.95,4.95,min(0.14*len(meta)+1.0,6),0.56),mix(tuple(WHITE),tuple(a),.18),meta,15,WHITE,bold=True)

# ---------------------------------------------------------------- GOALS (numbered list)
def s_goals(prs,sect,goals):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"学习目标 · learning goals","今天我们要学会什么?","Today's learning goals",a)
    area=Box(0.9,top+0.05,SW-1.8,SH-top-0.75)
    n=len(goals); gap=0.2
    gh=(area.h-gap*(n-1))/n
    for i,(zh,en) in enumerate(goals):
        b=Box(area.x,area.y+i*(gh+gap),area.w,gh)
        card(s,b)
        # number badge
        bs=min(gh-0.26,0.7)
        bb=Box(b.x+0.22,b.y+(gh-bs)/2,bs,bs)
        rect(s,bb,a,radius=0.24)
        text(s,bb,[[T(str(i+1),22,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        bilingual(s,Box(bb.x+bs+0.3,b.y,b.w-bs-1.0,gh),zh,en,zh_sz=18,en_sz=12)

# ---------------------------------------------------------------- 4-BIN OVERVIEW
def s_bins_overview(prs,sect):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"认识四个桶 · four bins","它们有什么不同?看颜色和符号","Four bins — tell them apart by color & symbol",a)
    area=Box(0.9,top+0.1,SW-1.8,SH-top-0.85)
    n=4; gap=0.26; cw=(area.w-gap*(n-1))/n
    for i,bn in enumerate(BINS):
        b=Box(area.x+i*(cw+gap),area.y,cw,area.h)
        card(s,b,shadow=True)
        # colored band carries the bin NAME (background color); frees the body for a big picture
        hb=Box(b.x,b.y,b.w,1.2)
        rect(s,hb,bn['c'],radius=0.11)
        rect(s,Box(b.x,b.y+0.75,b.w,0.45),bn['c'],radius=0)  # square off the bottom edge
        text(s,Box(b.x,b.y+0.16,b.w,1.0),
             [[T(bn['zh'],20,WHITE,True,False,DISP)],[T(bn['en'],11,mix(tuple(WHITE),tuple(bn['c']),.8),False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.1)
        # BIG photo placeholder — fills the card body (lots of room for a real photo)
        ph=Box(b.x+0.22,b.y+1.4,b.w-0.44,b.h-1.4-0.98)
        rect(s,ph,H('EAF0E2'),line=mix(tuple(bn['c']),(255,255,255),.4),line_w=1.5,radius=0.1,dash=True)
        text(s,ph,[[T("📷",40,mix(tuple(bn['c']),(255,255,255),.25))],
                   [T("图片 photo",12,mix(tuple(bn['c']),(255,255,255),.45),False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.05)
        rect(s,Box(b.x+0.2,b.y+b.h-0.92,b.w-0.4,0.74),SOFT,radius=0.16)
        text(s,Box(b.x+0.2,b.y+b.h-0.92,b.w-0.4,0.74),
             [[T(bn['dest'],11.5,GREEN,True)],[T(bn['desten'],9.5,GRAY,False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)

# ---------------------------------------------------------------- BIN DETAIL (identity + examples)
def s_bin_detail(prs,sect,idx):
    s=new_slide(prs); a=SECT[sect]['accent']; bn=BINS[idx]
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,f"认识 {idx+1}/4 · {bn['en'].lower()}",f"这是「{bn['zh']}」桶",f"Meet the {bn['en'].lower()} bin",a)
    area=Box(0.9,top+0.05,SW-1.8,SH-top-0.75)
    # left color identity panel
    lw=4.1
    lb=Box(area.x,area.y,lw,area.h)
    rect(s,lb,bn['c'],radius=0.06,shadow=True)
    # BIG identity photo placeholder (panel stays a background color)
    ph=Box(lb.x+0.35,lb.y+0.3,lw-0.7,area.h*0.52)
    rect(s,ph,WHITE,radius=0.12)
    text(s,ph,[[T("📷",48,bn['c'])],[T("图片 photo · 插入照片",12,GRAY,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.05)
    text(s,Box(lb.x,lb.y+area.h*0.52+0.42,lw,1.0),
         [[T(bn['zh'],26,WHITE,True,False,DISP)],
          [T(bn['py']+'  ·  '+bn['en'],12,mix(tuple(WHITE),tuple(bn['c']),.78),False,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.2)
    db=Box(lb.x+0.3,lb.y+lb.h-0.82,lb.w-0.6,0.66)
    rect(s,db,mix(tuple(WHITE),tuple(bn['c']),.9),radius=0.2)
    text(s,db,[[T(bn['dest'],12.5,bn['c'],True)],[T(bn['desten'],9.5,GRAY,False,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
    # right: examples grid
    rb=Box(lb.x+lw+0.35,area.y,area.w-lw-0.35,area.h)
    text(s,Box(rb.x,rb.y,rb.w,0.4),[[T("桶里有什么?  What goes in?",15,INK,True)]])
    grid_h=rb.h-0.5-(0.78 if bn.get('harm') else 0)
    grid=Box(rb.x,rb.y+0.5,rb.w,grid_h)
    ex=bn['ex']; cols=min(len(ex),3); rows=(len(ex)+cols-1)//cols
    gap=0.2; iw=(grid.w-gap*(cols-1))/cols; ih=min((grid.h-gap*(rows-1))/rows,2.1)
    for j,(em,zh,en) in enumerate(ex):
        r=j//cols; c=j%cols
        ib=Box(grid.x+c*(iw+gap),grid.y+r*(ih+gap),iw,ih)
        card(s,ib,line=mix(tuple(bn['c']),(255,255,255),.45),line_w=1.5,shadow=False)
        # BIGGER item picture placeholder (no emoji)
        ph2=Box(ib.x+0.12,ib.y+0.12,ib.w-0.24,ih*0.66)
        rect(s,ph2,H('EAF0E2'),line=mix(tuple(bn['c']),(255,255,255),.45),line_w=1.0,radius=0.08,dash=True)
        text(s,ph2,[[T("📷",26,mix(tuple(bn['c']),(255,255,255),.25))],
                    [T("图片 photo",9,mix(tuple(bn['c']),(255,255,255),.45),False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.0)
        text(s,Box(ib.x,ib.y+ih*0.66+0.1,ib.w,ih*0.34-0.1),
             [[T(zh,14,INK,True,False,DISP)],[T(en,10,GRAY,False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.05)
    if bn.get('harm'):
        safety_bar(s,Box(rb.x,rb.y+rb.h-0.66,rb.w,0.6),"不要碰!要告诉大人。","Don't touch — tell an adult.")
    return s

# ---------------------------------------------------------------- shared bits
def safety_bar(slide,box,zh,en):
    rect(slide,box,B_HARM,radius=0.16)
    text(slide,Box(box.x+0.2,box.y,box.w-0.4,box.h),
         [[T("🛑  "+zh,15,WHITE,True),T("   "+en,11,mix(tuple(WHITE),tuple(B_HARM),.82),False,True)]],
         anchor=MSO_ANCHOR.MIDDLE)

def teacher_tip(slide,box,zh):
    rect(slide,box,H('FFF3DD'),line=H('F1D08C'),line_w=1.0,radius=0.5)
    text(slide,box,[[T("👩‍🏫  "+zh,12.5,H('A85B12'),True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def timer_chip(slide,box,big,lab):
    rect(slide,box,H('FFF3DD'),line=H('F1D08C'),line_w=1.2,radius=0.16)
    text(slide,box,[[T(big,26,ORANGE,True,False,DISP)],[T(lab,11,ORANGE,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.0)

def media_box(slide,box,kind,label):
    rect(slide,box,H('11201A') if kind=='video' else H('EAF0E2'),
         line=None if kind=='video' else mix(tuple(MOSS),(255,255,255),.4),line_w=1.5,radius=0.05,dash=(kind!='video'))
    ic="▶" if kind=='video' else "📷"
    col=mix(tuple(MOSS),tuple(WHITE),.55) if kind=='video' else GRAY
    text(slide,box,[[T(ic,38,col)],[T(label,13,col,True)],
                    [T("插入"+("视频" if kind=='video' else "照片")+" · insert "+kind,10,mix(tuple(col),tuple(WHITE),.6) if kind=='video' else H('AEB3A6'),False,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.15)

def bin_chips_row(slide,box,gap=0.22):
    n=4; cw=(box.w-gap*(n-1))/n
    for i,bn in enumerate(BINS):
        b=Box(box.x+i*(cw+gap),box.y,cw,box.h)
        rect(slide,b,bn['c'],radius=0.5)
        text(slide,b,[[T(bn['ic']+" "+bn['zh'],14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- SENTENCE FRAME
def s_frame(prs,sect,tag,zh_html,en,note=""):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"句型 · sentence frame","会用句型说一说","Say it with the sentence frame",a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-1.0)
    b=Box(area.x,area.y+0.2,area.w,2.3)
    rect(s,b,GREEN,radius=0.05,shadow=True)
    pill(s,Box(b.x+0.4,b.y+0.4,1.9,0.55),mix(tuple(WHITE),tuple(GREEN),.2),tag,14,WHITE,bold=True)
    text(s,Box(b.x+0.4,b.y+1.1,b.w-0.8,1.0),
         [[T(zh_html,30,WHITE,True,False,DISP)]],anchor=MSO_ANCHOR.TOP)
    text(s,Box(b.x+0.4,b.y+1.75,b.w-0.8,0.5),[[T(en,15,mix(tuple(WHITE),tuple(MOSS),.8),False,True)]])
    if note:
        pill(s,Box(area.x,b.y+b.h+0.3,min(0.13*len(note)+1.0,7),0.55),SOFT,note,13,GREEN,bold=True)

# ---------------------------------------------------------------- prototype
def proto():
    prs=Presentation(); prs.slide_width=IN(SW); prs.slide_height=IN(SH)
    s_cover(prs)
    s_divider(prs,1,"第一节 · 打败垃圾怪!","Session 1 — Sort & defeat the Trash Monster","10:00–10:45  ·  11:00–11:45")
    s_goals(prs,1,[
        ("认识 4 种垃圾分类","Recognize the 4 bins — 可回收 · 厨余 · 有害 · 其他"),
        ("明白乱丢垃圾会破坏环境，所以要分类","Littering harms nature — so we sort"),
        ("知道每个桶的垃圾去哪儿","Recycle / compost / special handling / landfill"),
        ("会用句型「这是___，它要放进___桶」说一说","Use the sentence frame to describe items"),
    ])
    s_bins_overview(prs,1)
    s_bin_detail(prs,1,0)
    s_bin_detail(prs,1,2)
    s_frame(prs,1,"句型","这是 ____ 。它要放进 ____ 桶。","This is ____ . It goes in the ____ bin.","先指一指，再完整说一句")
    prs.save(OUT)
    print("proto saved",OUT,len(prs.slides._sldIdLst),"slides")

# ---------------------------------------------------------------- STEPS (process row)
def s_steps(prs,sect,kicker,zh,en,steps,note=None,timer=None,chips=False):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,kicker,zh,en,a)
    bottom=SH-0.75-(0.6 if note else 0)-(0.0 if not chips else 0.7)
    area=Box(0.9,top+0.25,SW-1.8,bottom-top-0.25)
    n=len(steps); aw=0.5; gap=0.0
    cw=(area.w-aw*(n-1))/n
    ch=min(area.h,2.7); cy=area.y+(area.h-ch)/2
    x=area.x
    for i,(em,zh2,en2) in enumerate(steps):
        b=Box(x,cy,cw,ch); card(s,b)
        rect(s,Box(b.x+cw/2-0.32,b.y+0.22,0.64,0.64),a,radius=0.24)
        text(s,Box(b.x,b.y+0.22,cw,0.64),[[T(str(i+1),22,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        text(s,Box(b.x,b.y+1.0,cw,ch-1.0),[[T(em,40,INK)],[T(zh2,17,INK,True,False,DISP)],[T(en2,11,GRAY,False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.15)
        x+=cw
        if i<n-1:
            text(s,Box(x,cy,aw,ch),[[T("→",24,a,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
            x+=aw
    if timer:
        timer_chip(s,Box(area.x,cy-0.0,1.7,0.0+0.0) if False else Box(SW-2.6,cy-0.95,1.7,0.85),timer[0],timer[1])
    if chips:
        bin_chips_row(s,Box(area.x+1.2,bottom+0.12,area.w-2.4,0.5))
    if note:
        teacher_tip(s,Box(area.x,SH-1.15,min(0.14*len(note)+1.0,9.5),0.5),note)

# ---------------------------------------------------------------- MEDIA + NOTES
def s_media(prs,sect,kicker,zh,en,kind,medialabel,cards,note=None,link=None):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.1,SW-1.8,SH-top-(1.15 if note else 0.7))
    lw=area.w*0.56
    mb=Box(area.x,area.y,lw,area.h)
    media_box(s,mb,kind,medialabel)
    if kind=='video' and link:
        # clickable play button + URL caption over the video box
        pbx=Box(mb.x+lw/2-1.5,mb.y+mb.h*0.60,3.0,0.6)
        sp=rect(s,pbx,ORANGE,radius=0.5,shadow=True)
        text(s,pbx,[[T("▶  点击播放 Open video",14,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        try: sp.click_action.hyperlink.address=link
        except Exception: pass
        lkb=Box(mb.x+0.3,mb.y+mb.h-0.55,lw-0.6,0.42)
        rect(s,lkb,WHITE,radius=0.18)
        text(s,lkb,[[T("🔗 "+link,10.5,GREEN,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rx=area.x+lw+0.35; rb=Box(rx,area.y,area.w-lw-0.35,area.h)
    n=len(cards); gap=0.22; chh=(rb.h-gap*(n-1))/n
    for i,(h,lines) in enumerate(cards):
        b=Box(rb.x,rb.y+i*(chh+gap),rb.w,chh); card(s,b)
        paras=[[T(h,15,a if a!=GREEN else GREEN,True,False,DISP)]]
        for ln in lines: paras.append([T(ln,13,INK,False)])
        text(s,b.pad(0.2),paras,anchor=MSO_ANCHOR.MIDDLE,spacing=1.25,space_after=3)
    if note:
        teacher_tip(s,Box(area.x,SH-1.1,min(0.14*len(note)+1.0,9.5),0.5),note)

# ---------------------------------------------------------------- KEY IDEA (Q cards + panel)
def s_keyidea(prs,sect,kicker,zh,en,qs,idea_zh,idea_en):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.15,SW-1.8,SH-top-0.8)
    lw=area.w*0.5
    n=len(qs); gap=0.24; qh=(area.h-gap*(n-1))/n
    for i,(q,qe) in enumerate(qs):
        b=Box(area.x,area.y+i*(qh+gap),lw,qh); card(s,b)
        rect(s,Box(b.x+0.18,b.y+qh/2-0.26,0.52,0.52),mix(tuple(a),(255,255,255),.16),radius=0.3)
        text(s,Box(b.x+0.18,b.y+qh/2-0.26,0.52,0.52),[[T(str(i+1),18,a,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        bilingual(s,Box(b.x+0.9,b.y,b.w-1.1,qh),q,qe,zh_sz=17,en_sz=11.5)
    pb=Box(area.x+lw+0.35,area.y,area.w-lw-0.35,area.h)
    rect(s,pb,GREEN,radius=0.06,shadow=True)
    text(s,Box(pb.x+0.35,pb.y+0.3,pb.w-0.7,0.5),[[T("💡 小结 Key idea",15,STAR,True)]])
    text(s,Box(pb.x+0.35,pb.y+0.95,pb.w-0.7,pb.h-1.4),
         [[T(idea_zh,22,WHITE,True,False,DISP)],[T(idea_en,13,mix(tuple(WHITE),tuple(MOSS),.8),False,True)]],
         anchor=MSO_ANCHOR.MIDDLE,spacing=1.3)

# ---------------------------------------------------------------- PRACTICE (item -> bin)
def s_practice(prs,sect,i,total,item,en,emoji,binidx,why,frame_zh,frame_en):
    s=new_slide(prs); a=SECT[sect]['accent']; bn=BINS[binidx]
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,f"放进哪个桶? {i}/{total} · which bin?",f"{emoji} {item}，放进哪个桶?",f"Which bin for a {en}?",a)
    area=Box(0.9,top+0.1,SW-1.8,SH-top-0.85)
    # left big item
    lw=3.7; lb=Box(area.x,area.y,lw,area.h)
    card(s,lb,shadow=True)
    text(s,lb,[[T(emoji,80,INK)],[T(item,24,INK,True,False,DISP)],[T(en,12,GRAY,False,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.15)
    # right: 4 mini bins, correct highlighted, + why + frame
    rb=Box(lb.x+lw+0.4,area.y,area.w-lw-0.4,area.h)
    grid_h=rb.h*0.5
    n=4; gap=0.2; cw=(rb.w-gap*(n-1))/n
    for j,b2 in enumerate(BINS):
        correct=(j==binidx)
        mb=Box(rb.x+j*(cw+gap),rb.y,cw,grid_h)
        if correct:
            rect(s,mb,b2['c'],radius=0.12,shadow=True)
            text(s,mb,[[T(b2['ic']+" ✓",26,WHITE)],[T(b2['zh'],14,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
        else:
            rect(s,mb,H('F2F1EA'),line=LINE,line_w=1.2,radius=0.12)
            text(s,mb,[[T(b2['ic'],22,H('C7C9C0'))],[T(b2['zh'],13,H('B7B9B0'),True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
    # why card
    wb=Box(rb.x,rb.y+grid_h+0.22,rb.w,0.7)
    card(s,wb,shadow=False)
    rect(s,Box(wb.x,wb.y+0.08,0.1,wb.h-0.16),bn['c'],radius=0.4)
    text(s,Box(wb.x+0.3,wb.y,wb.w-0.5,wb.h),[[T(f"→ {bn['zh']}",15,bn['c'],True),T(f"   {why}",14,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    # frame / safety
    fb=Box(rb.x,rb.y+grid_h+1.05,rb.w,rb.h-grid_h-1.05)
    if bn.get('harm'):
        safety_bar(s,Box(fb.x,fb.y+0.05,fb.w,min(fb.h-0.1,0.62)),"电池是有害垃圾,不要碰!","Batteries are hazardous — don't touch.")
    else:
        rect(s,Box(fb.x,fb.y+0.05,fb.w,min(fb.h-0.1,0.62)),GREEN,radius=0.16)
        text(s,Box(fb.x+0.25,fb.y+0.05,fb.w-0.5,min(fb.h-0.1,0.62)),
             [[T("说一说:  ",12,STAR,True),T(frame_zh,15,WHITE,True),T("  "+frame_en,10.5,mix(tuple(WHITE),tuple(MOSS),.8),False,True)]],anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- PRACTICE ANSWERS (reveal after the question slides)
def s_practice_answers(prs,sect,items):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"答案揭晓 · answers","放进哪个桶? — 一起看答案!","Which bin? — Now check the answers!",a)
    area=Box(0.9,top+0.15,SW-1.8,SH-top-0.8)
    n=len(items); gap=0.18; rh=(area.h-gap*(n-1))/n
    for k,(item,en,emoji,bi,why) in enumerate(items):
        bn=BINS[bi]
        b=Box(area.x,area.y+k*(rh+gap),area.w,rh); card(s,b,shadow=False)
        text(s,Box(b.x+0.35,b.y,3.2,rh),[[T(item,17,INK,True,False,DISP),T("  "+en,10,GRAY,False,True)]],anchor=MSO_ANCHOR.MIDDLE)
        text(s,Box(b.x+3.5,b.y,0.6,rh),[[T("→",20,a,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        pill(s,Box(b.x+4.15,b.y+rh/2-0.24,2.5,0.48),bn['c'],bn['zh']+"桶",14,WHITE,bold=True,fnt=DISP)
        paras=[[T(why,13,INK,False)]]
        if bn.get('harm'): paras.append([T("🛑 不要碰,告诉大人!",11,B_HARM,True)])
        text(s,Box(b.x+6.9,b.y,b.w-7.2,rh),paras,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
    return s

# ---------------------------------------------------------------- SORT GAME
def s_sortgame(prs):
    s=new_slide(prs); a=GREEN
    base(s,a); section_tab_pill(s,SECT[1]['tab'],a)
    top=header(s,"分类游戏 · 5 分钟 · sort game","打败垃圾怪 — 把垃圾送回家!","Sort it out — defeat the Trash Monster!",a)
    CARDS=[('📦','纸盒'),('🧴','塑料瓶'),('🍂','树叶'),('🍚','剩饭'),('🔋','旧电池'),('💡','旧灯泡'),('🧻','脏纸巾'),('🛍️','塑料袋')]
    area=Box(0.9,top+0.1,SW-1.8,SH-top-0.8)
    # card strip
    text(s,Box(area.x,area.y,area.w,0.35),[[T("🃏 要分类的垃圾卡  ·  the trash cards",13,INK,True)]])
    sb=Box(area.x,area.y+0.45,area.w,0.7)
    n=len(CARDS); gap=0.16; cw=(sb.w-gap*(n-1))/n
    for i,(em,zh) in enumerate(CARDS):
        b=Box(sb.x+i*(cw+gap),sb.y,cw,sb.h)
        rect(s,b,H('FFFFFF'),line=LINE,line_w=1.2,radius=0.16)
        text(s,b,[[T(em+" "+zh,12.5,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # 4 bins to drop into
    gb=Box(area.x,area.y+1.4,area.w,area.h-1.4)
    nn=4; g2=0.26; bw=(gb.w-g2*(nn-1))/nn
    for i,bn in enumerate(BINS):
        b=Box(gb.x+i*(bw+g2),gb.y,bw,gb.h)
        rect(s,b,bn['c'],radius=0.06,shadow=True)
        text(s,Box(b.x,b.y+0.3,b.w,1.1),[[T(bn['ic'],34,WHITE)],[T(bn['zh'],18,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
        db=Box(b.x+0.25,b.y+b.h-0.7,b.w-0.5,0.5)
        rect(s,db,mix(tuple(WHITE),tuple(bn['c']),.85),radius=0.5)
        text(s,db,[[T("把卡片放这里 drop here",10.5,bn['c'],True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- REVEAL + medals + answer key
def s_reveal(prs):
    s=new_slide(prs); a=GREEN
    base(s,a); section_tab_pill(s,SECT[1]['tab'],a)
    top=header(s,"揭晓 · reveal","垃圾怪被打败了! 🎉","Monster defeated — the town is clean!",a)
    area=Box(0.9,top+0.1,SW-1.8,SH-top-0.8)
    medals=[(H('F5B921'),H('5A4500'),'🥇','全对小队','All correct!'),
            (H('B9BCC2'),H('45474B'),'🥈','很棒小队','Almost!'),
            (H('CD7F32'),WHITE,'🥉','加油小队','Good try!')]
    mh=area.h*0.46; gap=0.3; mw=(area.w-gap*2)/3
    for i,(c,tc,em,zh,en) in enumerate(medals):
        b=Box(area.x+i*(mw+gap),area.y,mw,mh)
        rect(s,b,c,radius=0.08,shadow=True)
        text(s,b,[[T(em,34,tc)],[T(zh,18,tc,True,False,DISP)],[T(en,11,tc,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
    # answer key
    kb=Box(area.x,area.y+mh+0.25,area.w,area.h-mh-0.25)
    card(s,kb)
    text(s,Box(kb.x+0.3,kb.y+0.18,kb.w-0.6,0.4),[[T("✅ 答案 Answer key",14,GREEN,True)]])
    ANS=[('📦 纸盒',0),('🧴 塑料瓶',0),('🍂 树叶',1),('🍚 剩饭',1),('🔋 旧电池',2),('💡 旧灯泡',2),('🧻 脏纸巾',3),('🛍️ 塑料袋',3)]
    cols=4; cw=(kb.w-0.6)/cols
    for i,(lab,bi) in enumerate(ANS):
        r=i//cols; c=i%cols
        bx=kb.x+0.3+c*cw; by=kb.y+0.62+r*0.46
        text(s,Box(bx,by,cw,0.42),[[T(lab+"  ",13,INK,True),T("→ "+BINS[bi]['zh'],13,BINS[bi]['c'],True)]],anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- STATEMENT (frame + bins) exit/reflection/wrap
def s_statement(prs,sect,kicker,zh,en,frame_zh,frame_en,closing=None):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.3,SW-1.8,SH-top-1.0)
    fb=Box(area.x+0.6,area.y+0.1,area.w-1.2,1.9)
    rect(s,fb,PANEL,line=mix(tuple(a),(255,255,255),.4),line_w=2.5,radius=0.06,shadow=True)
    text(s,fb,[[T(frame_zh,30,GREEN,True,False,DISP)],[T(frame_en,14,GRAY,False,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.4)
    bin_chips_row(s,Box(area.x+1.4,fb.y+fb.h+0.4,area.w-2.8,0.55))
    if closing:
        pill(s,Box(area.x+area.w/2-3.0,area.y+area.h-0.0,6.0,0.0+0.0) if False else Box(SW/2-3.2,fb.y+fb.h+1.15,6.4,0.6),
             GREEN,closing,15,WHITE,bold=True,fnt=DISP)

# ---------------------------------------------------------------- VOCAB (我会认)
def s_vocab(prs,i,total,word):
    s=new_slide(prs); a=SECT[2]['accent']; c=word['c']
    base(s,a); section_tab_pill(s,SECT[2]['tab'],a)
    top=header(s,f"我会认 {i}/{total} · I can read",f"认一认:{word['ex']}","Recognize this word",a)
    area=Box(0.9,top+0.15,SW-1.8,SH-top-0.8)
    lw=area.w*0.46
    lb=Box(area.x,area.y,lw,area.h)
    card(s,lb,shadow=True,line=mix(tuple(c),(255,255,255),.45),line_w=2)
    text(s,lb,[[T(word['e'],58,INK)],[T(word['zh'],60,c,True,False,DISP)],
               [T(word['py'],20,GRAY,False,True)],[T(word['en'],17,INK,True)]],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.18)
    rb=Box(lb.x+lw+0.4,area.y,area.w-lw-0.4,area.h)
    # photo placeholder + usage word
    ph=Box(rb.x,rb.y,rb.w,rb.h*0.55)
    media_box(s,ph,'photo',f"{word['zh']} 实物照片")
    ub=Box(rb.x,rb.y+rb.h*0.55+0.25,rb.w,rb.h*0.45-0.25)
    card(s,ub)
    text(s,Box(ub.x+0.3,ub.y+0.2,ub.w-0.6,0.4),[[T("🗣️ 组词 Word",14,a,True)]])
    text(s,Box(ub.x+0.3,ub.y+0.6,ub.w-0.6,ub.h-0.8),[[T(word['ex'],30,c,True,False,DISP)]],anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- MATCHING
def s_match(prs):
    s=new_slide(prs); a=SECT[2]['accent']
    base(s,a); section_tab_pill(s,SECT[2]['tab'],a)
    top=header(s,"词汇配对 · 游戏 · matching","把桶的名字和符号连起来","Match each bin to its color symbol",a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-1.1)
    names=[b['zh'] for b in BINS]
    syms=[(BINS[2]['ic'],BINS[2]['c']),(BINS[0]['ic'],BINS[0]['c']),(BINS[3]['ic'],BINS[3]['c']),(BINS[1]['ic'],BINS[1]['c'])]
    n=4; gap=0.2; rh=(area.h-gap*(n-1))/n
    lw=area.w*0.42; rw=area.w*0.42
    for i in range(n):
        lb=Box(area.x,area.y+i*(rh+gap),lw,rh); card(s,lb)
        text(s,Box(lb.x+0.4,lb.y,lb.w-0.6,rh),[[T(names[i],22,BINS[i]['c'],True,False,DISP)]],anchor=MSO_ANCHOR.MIDDLE)
        rb=Box(area.x+area.w-rw,area.y+i*(rh+gap),rw,rh); card(s,rb)
        text(s,rb,[[T(syms[i][0],26,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(area.x+lw,area.y,area.w-lw-rw,area.h),[[T("⇔",30,a,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    teacher_tip(s,Box(area.x+area.w/2-2.2,SH-1.05,4.4,0.5),"请小朋友上来连线")

# ---------------------------------------------------------------- CARD GRID (slap cards / design ideas / take-home)
def s_cardgrid(prs,sect,kicker,zh,en,items,cols,note=None,big=44):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-(1.1 if note else 0.7))
    n=len(items); rows=(n+cols-1)//cols; gap=0.24
    cw=(area.w-gap*(cols-1))/cols; ch=(area.h-gap*(rows-1))/rows
    for i,(em,zh2,en2) in enumerate(items):
        r=i//cols; c=i%cols
        b=Box(area.x+c*(cw+gap),area.y+r*(ch+gap),cw,ch); card(s,b)
        paras=[[T(em,big,INK)]] if em else []
        paras.append([T(zh2,18,INK,True,False,DISP)])
        if en2: paras.append([T(en2,12,GRAY,False,True)])
        text(s,b.pad(0.12),paras,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.15)
    if note:
        teacher_tip(s,Box(area.x+area.w/2-3.0,SH-1.05,6.0,0.5),note)

# ---------------------------------------------------------------- DIALOG (pair share)
def s_dialog(prs):
    s=new_slide(prs); a=SECT[2]['accent']
    base(s,a); section_tab_pill(s,SECT[2]['tab'],a)
    top=header(s,"句型练习 · pair share","一起说一说","Practice the sentence frame",a)
    area=Box(0.9,top+0.15,SW-1.8,SH-top-0.8)
    lw=area.w*0.46
    # two frame cards
    f=[("这是 ____ 。","This is ____ ."),("它要放进 ____ 桶。","It goes in the ____ bin.")]
    fh=(area.h-0.24)/2
    for i,(zh,en) in enumerate(f):
        b=Box(area.x,area.y+i*(fh+0.24),lw,fh); card(s,b)
        rect(s,Box(b.x,b.y+0.1,0.1,b.h-0.2),(MOSS if i else GREEN),radius=0.4)
        text(s,Box(b.x+0.35,b.y,b.w-0.5,b.h),[[T(zh,23,INK,True,False,DISP)],[T(en,12.5,GRAY,False,True)]],anchor=MSO_ANCHOR.MIDDLE,spacing=1.2)
    # dialog example
    rb=Box(area.x+lw+0.4,area.y,area.w-lw-0.4,area.h); card(s,rb)
    text(s,Box(rb.x+0.35,rb.y+0.25,rb.w-0.7,0.45),[[T("👫 对话示范 Example",15,a,True)]])
    lines=[("A:","这是什么?",B_RECYCLE),("B:","这是香蕉皮。",B_FOOD),("A:","它要放进哪个桶?",B_RECYCLE),("B:","它要放进厨余垃圾桶!",B_FOOD)]
    y=rb.y+0.85
    for who,txt,col in lines:
        text(s,Box(rb.x+0.4,y,rb.w-0.8,0.5),[[T(who+" ",16,col,True),T(txt,16,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
        y+=0.62

# ---------------------------------------------------------------- WRITE (田字格)
def s_write(prs,i,total,zh,py,en,strokes):
    # day4_emergency-style: big-char card LEFT + 3-step panel + 田字格 + teacher/student bar
    s=new_slide(prs); a=SECT[2]['accent']
    base(s,a); section_tab_pill(s,SECT[2]['tab'],a)
    top=header(s,f"我会写 {i}/{total} · I can write",f"写一写:{zh}",f"I can write — {en}",a)
    area=Box(0.9,top+0.15,SW-1.8,SH-top-1.05)
    WARM=H('FFF6E6')
    # LEFT — big-character card
    lw=area.w*0.46; lb=Box(area.x,area.y,lw,area.h)
    rect(s,lb,WARM,line=a,line_w=2.5,radius=0.05,shadow=True)
    csz=130 if len(zh)==1 else 88
    text(s,Box(lb.x,lb.y+0.25,lb.w,area.h*0.55),[[T(zh,csz,a,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(lb.x,lb.y+area.h*0.64,lb.w,0.5),[[T(f"{py}  ·  {en}",20,GRAY,True,False,DISP)]],align=PP_ALIGN.CENTER)
    text(s,Box(lb.x,lb.y+area.h*0.64+0.55,lb.w,0.45),[[T(f"{strokes} 笔 / strokes",16,a,True,False,DISP)]],align=PP_ALIGN.CENTER)
    # RIGHT — 3-step panel
    rx=lb.x+lw+0.4; rw=area.w-lw-0.4
    sb=Box(rx,area.y,rw,1.75); card(s,sb)
    text(s,Box(sb.x+0.32,sb.y+0.2,rw-0.6,0.4),[[T("✍️ 3 步练习  3 Steps",16,a if a!=GREEN else GREEN,True,False,DISP)]])
    for k,st in enumerate(["1️⃣ 看老师写  Watch teacher","2️⃣ 用手指空中写  Air-write","3️⃣ 在田字格写 3 次  Write 3×"]):
        text(s,Box(sb.x+0.45,sb.y+0.66+k*0.35,rw-0.8,0.33),[[T(st,13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    # RIGHT — 4 田字格 cells
    gy=area.y+1.95; side=min((rw-0.28*3)/4,1.35)
    gx0=rx+(rw-(side*4+0.28*3))/2
    for k in range(4):
        bx=gx0+k*(side+0.28)
        rect(s,Box(bx,gy,side,side),WHITE,line=a,line_w=1.8,radius=0.0)
        cx=bx+side/2; cy=gy+side/2
        g1=s.shapes.add_connector(2,IN(cx),IN(gy),IN(cx),IN(gy+side)); g1.line.color.rgb=mix(tuple(a),(255,255,255),.4); g1.line.width=Pt(0.75); g1.line.dash_style=2
        g2=s.shapes.add_connector(2,IN(bx),IN(cy),IN(bx+side),IN(cy)); g2.line.color.rgb=mix(tuple(a),(255,255,255),.4); g2.line.width=Pt(0.75); g2.line.dash_style=2
    text(s,Box(rx,gy+side+0.12,rw,0.3),[[T("在田字格里写 3 次 ↓",12,GRAY,False)]],align=PP_ALIGN.CENTER)
    # BOTTOM — teacher asks / student does
    bb=Box(area.x,SH-0.98,area.w,0.62); rect(s,bb,WARM,line=a,line_w=1.5,radius=0.16)
    text(s,Box(bb.x+0.3,bb.y,bb.w*0.5-0.35,bb.h),[[T("👩‍🏫 老师问 Teacher asks: ",11,a,True),T(f"和我一起写「{zh}」",12,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(bb.x+bb.w*0.5+0.1,bb.y,bb.w*0.5-0.35,bb.h),[[T("🧒 学生 Student does: ",11,a,True),T("看 → 空中写 → 田字格写 3 次",12,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    return s

# ---------------------------------------------------------------- BAMBOOZLE
def s_bamboozle(prs):
    s=new_slide(prs); a=SECT[2]['accent']
    base(s,a); section_tab_pill(s,SECT[2]['tab'],a)
    top=header(s,"游戏时间 · Bamboozle","🎮 垃圾分类大闯关!","Team quiz game time",a)
    area=Box(0.9,top+0.15,SW-1.8,SH-top-0.8)
    lw=area.w*0.42
    lb=Box(area.x,area.y,lw,area.h)
    rect(s,lb,H('FFF3DD'),line=H('F1D08C'),line_w=1.4,radius=0.06,shadow=True)
    text(s,Box(lb.x,lb.y+0.5,lb.w,1.2),[[T("🪣 ♻️",46,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bb=Box(lb.x+0.6,lb.y+1.9,lb.w-1.2,0.8)
    rect(s,bb,ORANGE,radius=0.18,shadow=True)
    text(s,bb,[[T("▶  打开 Bamboozle",18,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(lb.x+0.3,lb.y+lb.h-0.7,lb.w-0.6,0.5),[[T("baamboozle.com — 老师课前创建游戏",11.5,H('A85B12'),True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rb=Box(lb.x+lw+0.4,area.y,area.w-lw-0.4,area.h)
    tb=Box(rb.x,rb.y,rb.w,rb.h*0.72); card(s,tb)
    text(s,Box(tb.x+0.35,tb.y+0.28,tb.w-0.7,0.5),[[T("📋 老师准备 Setup",15,a,True)]])
    steps=["1. 登录 baamboozle.com","2. 导入题库 bamboozle_day1_trash.csv","3. 分 2–3 组,轮流抢答","4. 「这个放哪个桶?」答对加分!"]
    y=tb.y+0.85
    for st in steps:
        text(s,Box(tb.x+0.45,y,tb.w-0.8,0.45),[[T(st,14,INK,False)]],anchor=MSO_ANCHOR.MIDDLE); y+=0.52
    pill(s,Box(rb.x,rb.y+rb.h*0.72+0.2,rb.w,rb.h*0.28-0.2),SOFT,"📂 约 10 题,覆盖 4 个桶",13,GREEN,bold=True)

# ---------------------------------------------------------------- CRAFT MENU
def s_craftmenu(prs):
    s=new_slide(prs); a=ORANGE
    base(s,a); section_tab_pill(s,SECT[3]['tab'],a)
    top=header(s,"选一选 · 做两个 · projects","今天做两个环保作品!","Two recycled-craft projects today",a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-0.8)
    projs=[(B_FOOD,'🎡','垃圾分类转盘','Project 1 · Sorting Spinner — 4 bins','纸盘 · 图钉 · 彩笔'),
           (ORANGE,'🔑','环保钥匙牌','Project 2 · Shrinky Dink Keychain','收缩片 · 马克笔 · 钥匙环')]
    gap=0.4; cw=(area.w-gap)/2
    for i,(c,em,zh,en,mat) in enumerate(projs):
        b=Box(area.x+i*(cw+gap),area.y,cw,area.h)
        card(s,b,shadow=True,line=mix(tuple(c),(255,255,255),.4),line_w=2)
        text(s,Box(b.x,b.y+0.6,b.w,1.3),[[T(em,64,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        text(s,Box(b.x,b.y+2.0,b.w,1.0),[[T(zh,28,c,True,False,DISP)],[T(en,13,GRAY,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.25)
        pill(s,Box(b.x+b.w/2-1.8,b.y+b.h-0.85,3.6,0.55),SOFT,mat,13,GREEN,bold=True,fnt=DISP)

# ---------------------------------------------------------------- MATERIALS
def s_materials(prs,kicker,zh,en,items):
    s=new_slide(prs); a=ORANGE
    base(s,a); section_tab_pill(s,SECT[3]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.3,SW-1.8,SH-top-0.9)
    n=len(items); gap=0.3; cw=(area.w-gap*(n-1))/n; ch=min(area.h,2.6); cy=area.y+(area.h-ch)/2
    for i,(em,zh2,en2) in enumerate(items):
        b=Box(area.x+i*(cw+gap),cy,cw,ch); card(s,b,shadow=True)
        text(s,b.pad(0.15),[[T(em,52,INK)],[T(zh2,19,INK,True,False,DISP)],[T(en2,11.5,GRAY,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.2)

print("builders loaded")

# ---------------------------------------------------------------- WORK TIME (timer + prompts)
def s_worktime(prs,kicker,zh,en,timer,head,prompts,danger=False):
    s=new_slide(prs); a=ORANGE
    base(s,a); section_tab_pill(s,SECT[3]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-0.8)
    tb=Box(area.x,area.y+area.h/2-0.9,2.2,1.8)
    timer_chip(s,tb,timer[0],timer[1])
    rb=Box(tb.x+2.2+0.4,area.y,area.w-2.2-0.4,area.h)
    hc=B_HARM if danger else a
    card(s,rb,line=(mix(tuple(B_HARM),(255,255,255),.0) if danger else LINE),line_w=(2.5 if danger else 1.0),
         fill=(H('FFF5F5') if danger else PANEL))
    text(s,Box(rb.x+0.4,rb.y+0.3,rb.w-0.8,0.5),[[T(head,16,hc,True,False,DISP)]])
    y=rb.y+0.95
    for p in prompts:
        text(s,Box(rb.x+0.5,y,rb.w-1.0,0.55),[[T(p,15,INK,False)]],anchor=MSO_ANCHOR.MIDDLE); y+=0.62

# ---------------------------------------------------------------- SHARE frame
def s_share(prs,kicker,zh,en,tag,fz,fe,emoji,sub,accent=ORANGE):
    s=new_slide(prs); a=accent
    base(s,a); section_tab_pill(s,SECT[3]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.3,SW-1.8,SH-top-1.0)
    fb=Box(area.x,area.y+0.1,area.w,1.7)
    rect(s,fb,GREEN if accent==B_FOOD else accent,radius=0.06,shadow=True)
    pill(s,Box(fb.x+0.4,fb.y+fb.h/2-0.28,2.0,0.56),mix(tuple(WHITE),tuple(GREEN if accent==B_FOOD else accent),.2),tag,14,WHITE,bold=True,fnt=DISP)
    text(s,Box(fb.x+2.7,fb.y,fb.w-3.0,fb.h),[[T(fz,26,WHITE,True,False,DISP)],[T(fe,13,mix(tuple(WHITE),tuple(MOSS),.8),False,True)]],anchor=MSO_ANCHOR.MIDDLE,spacing=1.3)
    cb=Box(area.x+area.w/2-3.2,fb.y+fb.h+0.35,6.4,area.h-fb.h-0.45)
    card(s,cb)
    text(s,cb,[[T(emoji,46,INK)],[T(sub,18,INK,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.2)

# ---------------------------------------------------------------- GROUP PHOTO
def s_groupphoto(prs):
    s=new_slide(prs); a=MOSS
    rect(s,Box(0,0,SW,SH),GREEN,radius=0)
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,IN(9.0),IN(-2.0),IN(7),IN(7))
    c.fill.solid(); c.fill.fore_color.rgb=mix(tuple(MOSS),tuple(GREEN),.5); c.line.fill.background(); c.shadow.inherit=False
    text(s,Box(0,2.0,SW,1.2),[[T("📸",60,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(0,3.2,SW,1.0),[[T("全班合影 — 拿着作品!",40,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(0,4.35,SW,0.5),[[T("Group photo — hold up your work! 🎉",17,mix(tuple(WHITE),tuple(MOSS),.85),False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Box(0,5.1,SW,0.8),[[T("🎡   🔑   ♻️   🌍   🌱",36,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- THREE CARDS (take-home)
def s_threecards(prs,sect,kicker,zh,en,cards):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,kicker,zh,en,a)
    area=Box(0.9,top+0.3,SW-1.8,SH-top-0.9)
    n=len(cards); gap=0.35; cw=(area.w-gap*(n-1))/n
    for i,(em,head,lines) in enumerate(cards):
        b=Box(area.x+i*(cw+gap),area.y,cw,area.h); card(s,b,shadow=True)
        text(s,Box(b.x,b.y+0.4,b.w,1.0),[[T(em,46,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        text(s,Box(b.x,b.y+1.5,b.w,0.5),[[T(head,17,a,True,False,DISP)]],align=PP_ALIGN.CENTER)
        paras=[[T(ln,13.5,INK,False)] for ln in lines]
        text(s,Box(b.x+0.3,b.y+2.1,b.w-0.6,b.h-2.3),paras,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.35,space_after=4)

# ---------------------------------------------------------------- INTERACTIVE: stand-by-your-bin (full-body vote)
def s_activity_vote(prs,sect):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"动一动 · 全身游戏 · move & sort","🏃 站到你的桶!","Run to the right bin — sorting race",a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-1.15)
    corners=["教室 左前角","教室 右前角","教室 左后角","教室 右后角"]
    n=4; gap=0.26; cw=(area.w-gap*(n-1))/n
    for i,bn in enumerate(BINS):
        b=Box(area.x+i*(cw+gap),area.y,cw,area.h); rect(s,b,bn['c'],radius=0.07,shadow=True)
        text(s,Box(b.x,b.y+0.35,b.w,1.2),[[T(bn['ic'],34,WHITE)],[T(bn['zh'],17,WHITE,True,False,DISP)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spacing=1.1)
        db=Box(b.x+0.22,b.y+b.h-0.72,b.w-0.44,0.52); rect(s,db,mix(tuple(WHITE),tuple(bn['c']),.85),radius=0.5)
        text(s,db,[[T("📍 "+corners[i],10.5,bn['c'],True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    teacher_tip(s,Box(area.x+area.w/2-4.4,SH-1.05,8.8,0.5),
                "老师举垃圾卡 → 学生跑到对应桶的角落 → 一起说「这是___,放进___桶!」")

# ---------------------------------------------------------------- INTERACTIVE: action chant (TPR)
def s_activity_tpr(prs,sect):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"动一动 · 律动 · chant & move","🎵 垃圾分类律动操","Sorting action chant — do it together",a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-1.15)
    acts=[(BINS[0],"转一转手臂","spin your arms"),(BINS[1],"摸摸小肚子","rub your tummy"),
          (BINS[2],"双手交叉 — 停!","cross arms — stop!"),(BINS[3],"往后一推","push it away")]
    n=4; gap=0.26; cw=(area.w-gap*(n-1))/n
    for i,(bn,act,acten) in enumerate(acts):
        b=Box(area.x+i*(cw+gap),area.y,cw,area.h)
        card(s,b,shadow=True,line=mix(tuple(bn['c']),(255,255,255),.4),line_w=2)
        text(s,Box(b.x,b.y+0.35,b.w,1.0),[[T(bn['ic'],40,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        text(s,Box(b.x,b.y+1.5,b.w,0.5),[[T(bn['zh'],17,bn['c'],True,False,DISP)]],align=PP_ALIGN.CENTER)
        text(s,Box(b.x+0.12,b.y+2.05,b.w-0.24,b.h-2.2),
             [[T(act,14,INK,True,False,DISP)],[T(acten,10.5,GRAY,False,True)]],
             align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.15)
    pill(s,Box(area.x+area.w/2-3.6,SH-1.05,7.2,0.55),SOFT,
         "🗣️ 边做边说:「可回收 — 转! 厨余 — 摸! 有害 — 停! 其他 — 推!」",12.5,GREEN,bold=True,fnt=DISP)

# ---------------------------------------------------------------- REVIEW: Bamboozle (blank link placeholder)
def s_review_bamboozle(prs,sect,topic_label):
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"复习 · review game","🎮 Bamboozle 复习大闯关",f"Review game — {topic_label}",a)
    area=Box(0.9,top+0.2,SW-1.8,SH-top-0.8)
    lw=area.w*0.46
    lb=Box(area.x,area.y,lw,area.h)
    rect(s,lb,H('FFF3DD'),line=H('F1D08C'),line_w=1.4,radius=0.06,shadow=True)
    text(s,Box(lb.x,lb.y+0.45,lb.w,1.2),[[T("🎮 ♻️",48,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bb=Box(lb.x+lw/2-1.7,lb.y+1.95,3.4,0.8); rect(s,bb,ORANGE,radius=0.18,shadow=True)
    text(s,bb,[[T("▶  开始复习游戏",18,WHITE,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # blank link placeholder line
    lk=Box(lb.x+0.4,lb.y+lb.h-0.78,lw-0.8,0.55); rect(s,lk,WHITE,line=H('F1D08C'),line_w=1.2,radius=0.16)
    text(s,Box(lk.x+0.18,lk.y,lk.w-0.3,lk.h),
         [[T("🔗 链接 Link: ",12,H('A85B12'),True),T("________________  (老师课前粘贴)",12,GRAY,False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    rb=Box(lb.x+lw+0.4,area.y,area.w-lw-0.4,area.h); card(s,rb)
    text(s,Box(rb.x+0.35,rb.y+0.3,rb.w-0.7,0.5),[[T("📋 怎么玩 How to play",15,a if a!=GREEN else GREEN,True)]])
    steps=["1. 老师点开上面的链接","2. 分 2–3 组,轮流抢答","3. 答对加分,答错全班一起说正确答案",f"4. 复习今天学的:{topic_label}"]
    y=rb.y+0.95
    for st in steps:
        text(s,Box(rb.x+0.5,y,rb.w-1.0,0.5),[[T(st,14,INK,False)]],anchor=MSO_ANCHOR.MIDDLE); y+=0.62
    return s

# ---------------------------------------------------------------- COMPLETE BOOKLET (before projects)
def s_booklet(prs,sect,days):
    """days: list of (label, zh_topic) booklets to finish first."""
    s=new_slide(prs); a=SECT[sect]['accent']
    base(s,a); section_tab_pill(s,SECT[sect]['tab'],a)
    top=header(s,"先写 · 再动手 · booklet first","📓 先完成练习册!","Finish your booklet first — then crafts",a)
    area=Box(0.9,top+0.25,SW-1.8,SH-top-1.15)
    n=len(days); gap=0.35; cw=min((area.w-gap*(n-1))/n,3.6) if n>0 else area.w
    total=cw*n+gap*(n-1); x0=area.x+(area.w-total)/2
    for i,(lab,zh) in enumerate(days):
        b=Box(x0+i*(cw+gap),area.y,cw,area.h); card(s,b,shadow=True,line=mix(tuple(a),(255,255,255),.4),line_w=2)
        text(s,Box(b.x,b.y+0.45,b.w,1.0),[[T("📓",54,INK)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        text(s,Box(b.x,b.y+1.6,b.w,0.9),[[T(lab+" 练习册",20,a if a!=GREEN else GREEN,True,False,DISP)],
             [T(zh,13,GRAY,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,spacing=1.2)
        ck=Box(b.x+b.w/2-1.2,b.y+b.h-0.8,2.4,0.55); rect(s,ck,SOFT,radius=0.5)
        text(s,ck,[[T("✅ 完成 → 打勾",12.5,GREEN,True,False,DISP)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    teacher_tip(s,Box(area.x+area.w/2-4.0,SH-1.05,8.0,0.5),
                "做手工之前,先一起完成今天的练习册 — 先看 → 一起读 → 自己写 → 同桌检查")
    return s

# ================================================================ FULL DECK
RECOG=[
 {'e':'♻️','zh':'可回收','py':'kě huí shōu','en':'recyclable','ex':'可回收物','c':B_RECYCLE},
 {'e':'🍎','zh':'厨余','py':'chú yú','en':'food waste','ex':'厨余垃圾','c':B_FOOD},
 {'e':'⚠️','zh':'有害','py':'yǒu hài','en':'hazardous','ex':'有害垃圾','c':B_HARM},
 {'e':'🗑️','zh':'其他','py':'qí tā','en':'other','ex':'其他垃圾','c':B_OTHER},
 {'e':'🔀','zh':'分类','py':'fēn lèi','en':'to sort','ex':'垃圾分类','c':GREEN},
]
WRITE=[('垃圾','lā jī','trash','8 + 6'),('分类','fēn lèi','to sort','4 + 9'),('回收','huí shōu','recycle','6 + 6')]
SCEN=[
 ('旧书本','old book','📚',0,'纸做的、能再利用'),
 ('骨头','bone','🦴',1,'吃剩的、会腐烂'),
 ('废旧电池','old battery','🔋',2,'有毒,对环境有害'),
 ('脏纸巾','dirty tissue','🧻',3,'用过了、不能回收'),
 ('玻璃罐','glass jar','🫙',0,'洗干净能再用'),
]

def build():
    PAGE[0]=0
    prs=Presentation(); prs.slide_width=IN(SW); prs.slide_height=IN(SH)
    # ---- SESSION 1
    s_cover(prs)
    s_divider(prs,1,"垃圾哪里去了?","Session 1 — Where does our trash go?","10:00–10:45   ·   11:00–11:45")
    s_goals(prs,1,[
        ("认识 4 种垃圾分类","Recognize the 4 bins — 可回收 · 厨余 · 有害 · 其他"),
        ("明白乱丢垃圾会破坏环境,所以要分类","Littering harms nature — so we must sort"),
        ("知道每个桶的垃圾去哪儿","Recycle / compost / special handling / landfill"),
        ("会用句型「这是___,它要放进___桶」说一说","Use the sentence frame to describe items"),
    ])
    s_media(prs,1,"热身 · 看视频 · video hook","📺 垃圾哪里去了?","Where does our trash go?",
        'video',"垃圾哪里去了?  (点击播放)",
        [("👀 看之前 — 想一想 Before",["• 你扔的垃圾,最后去了哪里?","• 猜一猜:垃圾会自己消失吗?"]),
         ("🎯 看的时候 — 找一找 During",["• 垃圾被运到了哪几个地方?","• 哪些垃圾可以变成新东西?"])],
        note="老师:点击播放,边看边停 — 看完一起讨论(下一页)",
        link="https://www.youtube.com/watch?v=AKoKVQb9_80")
    s_keyidea(prs,1,"讨论 · 看完后 · after watching","垃圾到底去哪儿了?","Where did the trash really go?",
        [("垃圾去了哪几个地方?","填埋场 · 焚烧厂 · 回收站 — 大部分被埋起来 / 烧掉"),
         ("不分类会怎么样?","好垃圾和坏垃圾混在一起,不能回收,还会污染环境")],
        "垃圾不会消失!乱丢、不分类会破坏地球。我们要学会垃圾分类!","Trash doesn't vanish — sorting protects our planet!")
    s_steps(prs,1,"热身 · 3 分钟 · think-pair-share","我们能帮忙吗?想 · 说 · 听","Can we help? Think · Pair · Share",
        [("🤔","想一想","Think"),("👫","说一说","Pair"),("👂","听一听","Share")],
        timer=("3:00","⏱ 讨论"))
    s_bins_overview(prs,1)
    s_activity_vote(prs,1)                 # interactive: run to the right bin
    for i in range(4): s_bin_detail(prs,1,i)
    s_activity_tpr(prs,1)                  # interactive: sorting action chant
    s_frame(prs,1,"句型","这是 ____ 。它要放进 ____ 桶。","This is ____ . It goes in the ____ bin.","先指一指,再完整说一句")
    for k,(item,en,em,bi,why) in enumerate(SCEN):
        s_practice(prs,1,k+1,len(SCEN),item,en,em,bi,why,
                   f"这是{item}。它要放进{BINS[bi]['zh']}桶。",f"This is a {en}. It goes in the {BINS[bi]['en'].lower()} bin.")
    s_practice_answers(prs,1,SCEN)
    s_sortgame(prs)
    s_reveal(prs)
    s_statement(prs,1,"出门票 · exit ticket","离开前,说一句!","Say one sentence before you go",
        "这是 ____ 。它要放进 ____ 桶。","This is ____ . It goes in the ____ bin.")
    # ---- SESSION 2
    s_divider(prs,2,"第二节 · 词语与游戏","Session 2 — Words & Games","45 min")
    for i,w in enumerate(RECOG): s_vocab(prs,i+1,len(RECOG),w)
    s_match(prs)
    s_cardgrid(prs,2,"拍词卡 · 游戏 · slap cards","听到词,快拍它!","Slap the card you hear!",
        [("♻️","可回收",""),("🍎","厨余",""),("⚠️","有害",""),("🗑️","其他",""),("🔀","分类",""),("🪣","垃圾桶","")],
        cols=3,note="老师说词 → 最快拍中得 1 分")
    s_dialog(prs)
    for i,(zh,py,en,st) in enumerate(WRITE): s_write(prs,i+1,len(WRITE),zh,py,en,st)
    s_bamboozle(prs)
    s_review_bamboozle(prs,2,"垃圾分类 · 4 个桶")
    s_statement(prs,2,"小结 · reflection","今天我学会了……","Today I learned…",
        "今天我学会了 __________","Today I learned __________ .")
    # ---- SESSION 3
    s_divider(prs,3,"第三节 · 我是垃圾分类专家!","Session 3 — Two Recycled Crafts","90 min")
    s_booklet(prs,3,[("Day 1","垃圾分类 Sorting")])
    s_craftmenu(prs)
    s_materials(prs,"作品 1 · 材料 · materials","垃圾分类转盘 — 材料","Sorting spinner — materials",
        [("🍽️","纸盘","paper plate"),("📌","图钉","brass fastener"),("🖍️","彩笔","markers"),("✂️","剪刀","scissors")])
    s_steps(prs,3,"作品 1 · 做法 · steps","4 步做转盘","Spinner — 4 steps",
        [("🍽️","拿纸盘","Take a plate"),("➗","分 4 格","4 sections"),("🎨","画 4 个桶","Draw 4 bins"),("📌","装指针","Add pointer")],chips=True)
    s_media(prs,3,"作品 1 · 老师示范 · demo","看老师做好的转盘","Teacher demo — the finished spinner",
        'photo',"做好的转盘照片",
        [("🎡 怎么玩?",["转一转指针 →","指到哪个桶,就说:","「这是___。它要放进___桶!」"])])
    s_worktime(prs,"作品 1 · 动手做 · work time","开始做!老师巡视","Work time + teacher walk-around",
        ("30:00","⏱ 制作时间"),"🗣️ 巡视提问 Prompts",
        ["• 这个桶是什么颜色?","• 香蕉皮放哪个桶?","• 旧电池放哪个桶?(🛑 有害!)","• 用句型说给我听!"])
    s_materials(prs,"作品 2 · 材料 · materials","环保钥匙牌 — 材料","Shrinky Dink keychain — materials",
        [("🟦","收缩片","Shrinky Dink"),("🖊️","马克笔","Sharpies"),("🔑","钥匙环","keyring"),("🔥","烤箱","oven · 老师")])
    s_steps(prs,3,"作品 2 · 做法 · steps","4 步做钥匙牌","Keychain — 4 steps",
        [("🎨","画图案","Draw"),("✂️","剪下来","Cut out"),("🔥","烤一烤","Bake (teacher)"),("🔑","装钥匙环","Add keyring")])
    s_cardgrid(prs,3,"作品 2 · 设计灵感 · ideas","画什么好呢?","Design ideas",
        [("🦸‍♂️","回收小英雄","Recycle Hero"),("🌱⭐","环保小达人","Eco Star"),("🟦🟩🟥⬛","4 个桶图标","4 bin icons")],cols=3,big=50)
    s_worktime(prs,"作品 2 · 安全 + 动手 · safety","安全第一,开始做!","Safety first, then work time",
        ("25:00","⏱ 制作时间"),"⚠️ 烤箱安全 Oven safety",
        ["🔥 烤箱只有老师能碰 — 350°F,1–3 分钟","🙌 小朋友画好、剪好,交给老师烤","🧤 烤好的牌子要凉一凉再拿"],danger=True)
    s_steps(prs,3,"收尾 · gallery walk","🖼️ 作品展览会","Gallery Walk",
        [("🪑","摆桌","Place work"),("🚶","静走","Walk & look"),("👏","拍桌","Tap to applaud")],
        note="看到喜欢的作品,轻轻拍桌说「真棒!」")
    s_share(prs,"分享 1/2 · 转盘 · share","拿着转盘说一说","Share your spinner",
        "转盘 Spinner","我转到了 ____ 。它要放进 ____ 桶!","I landed on ____ . It goes in the ____ bin!","🎡","转一转,说给大家听!",accent=B_FOOD)
    s_share(prs,"分享 2/2 · 钥匙牌 · share","拿着钥匙牌说一说","Share your keychain",
        "钥匙牌 Keychain","我做了一个 ____ 。我是环保小达人!","I made a ____ . I am an Eco Hero!","🔑✨","挂在书包上,每天提醒自己分类!",accent=ORANGE)
    s_groupphoto(prs)
    # ---- WRAP UP
    s_threecards(prs,4,"活动延伸 · 在家练习 · extend","把垃圾分类带回家!","Take sorting home",
        [("🏫","在教室 Classroom",["设 2–4 个分类桶","每天自己分类丢垃圾"]),
         ("🏠","在家里 At home",["和爸爸妈妈一起分类","看动画《宝宝巴士·学垃圾分类》"]),
         ("🔬","科学区 Center",["垃圾卡片 + 垃圾桶底板","动手贴一贴、分一分"])])
    s_statement(prs,4,"总结 · wrap-up","今天我学会了……","What I learned today",
        "今天我学会了 ____ 要放进 ____ 桶。","Today I learned that ____ goes in the ____ bin.",
        closing="🌱 垃圾分类小专家 · 下次见!")
    prs.save(OUT)
    print("SAVED",OUT,len(prs.slides._sldIdLst),"slides")

if __name__=='__main__':
    build()
