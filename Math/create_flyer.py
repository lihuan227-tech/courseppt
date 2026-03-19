#!/usr/bin/env python3
"""Generate a one-page Math course flyer in PPT format with full session details."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- constants ---
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
DARK_ORANGE = RGBColor(0xE6, 0x7E, 0x00)
GREEN = RGBColor(0x9D, 0xC6, 0x4D)
DARK_GREEN = RGBColor(0x5E, 0x94, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xFF, 0xF8, 0xF0)
LIGHT_GREEN_BG = RGBColor(0xF1, 0xF8, 0xF3)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)

FONT = 'KaiTi'

prs = Presentation()
prs.slide_width = Inches(8.5)
prs.slide_height = Inches(11)

slide = prs.slides.add_slide(prs.slide_layouts[6])


def add_shape(left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=10, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = FONT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return p


def add_para(tf, text, font_size=10, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = Pt(0)
    p.font.name = FONT
    return p


# ============================
# HEADER BANNER (compact)
# ============================
banner = add_rect(Inches(0), Inches(0), Inches(8.5), Inches(0.75), fill_color=GREEN)
add_rect(Inches(0), Inches(0.72), Inches(8.5), Inches(0.03), fill_color=ORANGE)

txBox = add_textbox(Inches(0.5), Inches(0.08), Inches(7.5), Inches(0.38))
tf = txBox.text_frame
set_text(tf, '暑假数学课程  Summer Math Course', font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox = add_textbox(Inches(0.5), Inches(0.45), Inches(7.5), Inches(0.22))
tf = txBox.text_frame
set_text(tf, '谷雨中文 · 2025 Summer Program', font_size=10, color=RGBColor(0xFF, 0xFF, 0xE0), alignment=PP_ALIGN.CENTER)

# ============================
# FEATURES (single compact row)
# ============================
y_feat = Inches(0.82)
features = [
    ('🧮 系统化课程设计', '5个级别，由浅入深，涵盖计算、图形、应用题与逻辑推理'),
    ('🧩 思维训练为核心', '注重数学思维与解题策略培养，锻炼逻辑分析能力'),
    ('🚀 暑期衔接与提升', '巩固已学知识，预习新学期内容，查漏补缺、拓展提高'),
]

card_w = Inches(2.45)
card_h = Inches(0.62)
gap = Inches(0.12)
start_x = Inches(0.55)

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    add_shape(x, y_feat, card_w, card_h, fill_color=WHITE, border_color=ORANGE, border_width=Pt(1.5))

    txBox = add_textbox(x + Inches(0.08), y_feat + Inches(0.04), card_w - Inches(0.16), Inches(0.2))
    tf = txBox.text_frame
    set_text(tf, title, font_size=9, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    txBox = add_textbox(x + Inches(0.08), y_feat + Inches(0.25), card_w - Inches(0.16), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, desc, font_size=7, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ============================
# LEVELS WITH FULL SESSION DETAILS
# ============================

levels_data = [
    {
        'name': 'Level 1',
        'sessions': [
            {
                'name': 'S1',
                'lessons': [
                    'Add & Subtract Three-digit Numbers',
                    'Commutative & Associative Property',
                    'Addition Strategies',
                    'Word Problems',
                ],
                'focus': '三位数运算：加减三位数、运算律、加法策略与应用题',
            },
            {
                'name': 'S2',
                'lessons': [
                    'Subtraction Strategies',
                    'Addition & Subtraction Strategies',
                    'Find the Missing Numbers',
                    'Word Problems',
                ],
                'focus': '运算策略深化：减法策略、加减混合策略、找缺失数、应用题',
            },
            {
                'name': 'S3',
                'lessons': [
                    'Lining-up Problems',
                    'Age Puzzles',
                    'Two-Dimensional Shapes',
                    'Three-Dimensional Shapes',
                ],
                'focus': '逻辑推理与图形：排队问题、年龄谜题、平面与立体图形识别',
            },
            {
                'name': 'S4',
                'lessons': [
                    'Bar Graphs',
                    'Tally Charts',
                    'Directional Words',
                    'Guess the Number',
                ],
                'focus': '数据与方向：条形图、计数表、方向词、猜数游戏',
            },
        ],
    },
    {
        'name': 'Level 2',
        'sessions': [
            {
                'name': 'S1',
                'lessons': [
                    'Even and Odd Numbers I',
                    'Even and Odd Numbers II',
                    'Multiplication (Vertical Method)',
                    'Word Problems',
                ],
                'focus': '奇偶与乘法：奇偶数规律、竖式乘法、应用题',
            },
            {
                'name': 'S2',
                'lessons': [
                    'Parentheses I',
                    'Parentheses II',
                    'Distributive Property',
                    'Word Problems',
                ],
                'focus': '括号与分配律：括号运算、分配律及应用题',
            },
            {
                'name': 'S3',
                'lessons': [
                    'Multiply & Divide with 10, 100, 1000',
                    'Division with Remainders',
                    'Division with Distributive Property',
                    'Word Problems',
                ],
                'focus': '除法深化：乘除10/100/1000、余数除法、分配律除法',
            },
            {
                'name': 'S4',
                'lessons': [
                    'Magic Square',
                    'Page Number Word Problems',
                    'Time Word Problems',
                    'Project: Find & Draw the Pattern',
                ],
                'focus': '逻辑与规律：幻方、页码问题、时间应用题、寻找规律',
            },
        ],
    },
    {
        'name': 'Level 3',
        'sessions': [
            {
                'name': 'S1',
                'lessons': [
                    'Order of Operations',
                    'Parentheses I',
                    'Parentheses II',
                    'Word Problems',
                ],
                'focus': '运算顺序与括号：四则运算法则、括号运用、应用题',
            },
            {
                'name': 'S2',
                'lessons': [
                    'GCF and LCM',
                    'Mixed Numbers & Improper Fractions',
                    'Add & Subtract Fractions',
                    'Add & Subtract Mixed Numbers',
                ],
                'focus': '分数基础：GCF/LCM、带分数与假分数互换、分数加减',
            },
            {
                'name': 'S3',
                'lessons': [
                    'Multiply Fractions',
                    'Multiply Fractions by Whole Numbers',
                    'Fractions Word Problems I',
                    'Fractions Word Problems II',
                ],
                'focus': '分数乘法与应用：分数乘法、分数乘整数、应用题',
            },
            {
                'name': 'S4',
                'lessons': [
                    'Bar Graphs Word Problems I',
                    'Bar Graphs Word Problems II',
                    'Chicken and Rabbit Problems',
                    'Project: Find & Draw the Pattern',
                ],
                'focus': '图表与逻辑：条形图应用、鸡兔同笼、寻找规律',
            },
        ],
    },
    {
        'name': 'Level 4',
        'sessions': [
            {
                'name': 'S1',
                'lessons': [
                    'Prime & Composite Numbers 1',
                    'Prime & Composite Numbers 2',
                    'Prime Factorization 1',
                    'Prime Factorization 2',
                ],
                'focus': '质数与因数分解：质数合数识别、质因数分解',
            },
            {
                'name': 'S2',
                'lessons': [
                    'Add, Subtract & Multiply Decimals',
                    'Divide Decimals',
                    'Convert Fractions to Decimals',
                    'Compare & Order Fractions/Decimals',
                ],
                'focus': '小数运算：加减乘除小数、分数与小数互转及比较',
            },
            {
                'name': 'S3',
                'lessons': [
                    'Simplify Fractions',
                    'Calculation Trick - Fraction 1',
                    'Calculation Trick - Fraction 2',
                    'Fractions ×÷ Word Problems',
                ],
                'focus': '分数深化：化简分数、分数计算技巧、乘除应用题',
            },
            {
                'name': 'S4',
                'lessons': [
                    'Convert to Percents',
                    'Percentage Word Problems',
                    'Average Word Problems 1',
                    'Average Word Problems 2',
                ],
                'focus': '百分比与平均数：分数/小数/比率转百分比、百分比应用、平均数',
            },
        ],
    },
    {
        'name': 'Level 5',
        'sessions': [
            {
                'name': 'S1',
                'lessons': [
                    'Operations of Negative Numbers',
                    'Operations of Fractions & Decimals',
                    'Operations of Percents',
                    'Operations of Exponents',
                ],
                'focus': '数的运算基础：负数、分数小数、百分数、指数的四则运算',
            },
            {
                'name': 'S2',
                'lessons': [
                    'Real Numbers & Simplifying Expressions',
                    'Linear Equations (One Variable)',
                    'Linear Inequalities',
                    'Applications of Linear Equations',
                ],
                'focus': '代数核心：实数性质化简、一元线性方程与不等式及应用',
            },
            {
                'name': 'S3',
                'lessons': [
                    'Parallel Lines and Angles',
                    'Properties of Triangles',
                    'Area & Circumference of Circles',
                    'Area & Circumference of Sectors',
                ],
                'focus': '平面几何：平行线与角、三角形性质、圆的周长与面积、扇形',
            },
            {
                'name': 'S4',
                'lessons': [
                    'Volume of Prism & Cylinder',
                    'Volume of Pyramid & Cone',
                    'Surface Area of Prism & Cylinder',
                    'Surface Area of Pyramid & Cone',
                ],
                'focus': '立体几何：棱柱/圆柱/棱锥/圆锥的体积与表面积',
            },
        ],
    },
]

y_start = Inches(1.52)
level_h = Inches(1.72)
level_gap = Inches(0.08)
page_x = Inches(0.35)
page_w = Inches(7.8)

for lvl_idx, level in enumerate(levels_data):
    y_lvl = y_start + lvl_idx * (level_h + level_gap)

    # Level label bar
    label_h = Inches(0.22)
    add_rect(page_x, y_lvl, page_w, label_h, fill_color=ORANGE)
    txBox = add_textbox(page_x + Inches(0.1), y_lvl + Inches(0.01), page_w - Inches(0.2), label_h)
    tf = txBox.text_frame
    set_text(tf, level['name'] + ' Summer Math Curriculum', font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 4 session cards in a row
    sess_y = y_lvl + label_h + Inches(0.04)
    sess_w = (page_w - 3 * Inches(0.08)) / 4
    sess_h = level_h - label_h - Inches(0.04)

    for s_idx, sess in enumerate(level['sessions']):
        sx = page_x + s_idx * (sess_w + Inches(0.08))

        # Session card background
        add_shape(sx, sess_y, sess_w, sess_h, fill_color=WHITE, border_color=RGBColor(0xE0, 0xE0, 0xE0), border_width=Pt(0.75))

        # Session header bar
        sh_h = Inches(0.18)
        add_rect(sx, sess_y, sess_w, sh_h, fill_color=DARK_GREEN)
        txBox = add_textbox(sx + Inches(0.02), sess_y + Inches(0.01), sess_w - Inches(0.04), sh_h)
        tf = txBox.text_frame
        set_text(tf, sess['name'], font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Lesson list
        lesson_y = sess_y + sh_h + Inches(0.02)
        txBox = add_textbox(sx + Inches(0.06), lesson_y, sess_w - Inches(0.12), Inches(0.72))
        tf = txBox.text_frame
        tf.word_wrap = True
        for li, lesson in enumerate(sess['lessons']):
            lesson_text = f"L{s_idx*4 + li + 1}  {lesson}"
            if li == 0:
                set_text(tf, lesson_text, font_size=6, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)
            else:
                add_para(tf, lesson_text, font_size=6, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(1.5))

        # Focus box
        focus_y = sess_y + sess_h - Inches(0.38)
        add_shape(sx + Inches(0.04), focus_y, sess_w - Inches(0.08), Inches(0.34),
                  fill_color=LIGHT_GREEN_BG)
        txBox = add_textbox(sx + Inches(0.08), focus_y + Inches(0.02), sess_w - Inches(0.16), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        set_text(tf, '📌 ' + sess['focus'], font_size=5.5, color=DARK_GREEN, alignment=PP_ALIGN.LEFT)

# ============================
# FOOTER
# ============================
y_footer = Inches(10.55)
add_rect(Inches(0), y_footer, Inches(8.5), Inches(0.45), fill_color=ORANGE)

txBox = add_textbox(Inches(0.5), y_footer + Inches(0.03), Inches(7.5), Inches(0.2))
tf = txBox.text_frame
set_text(tf, '立即报名  Register Now', font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

txBox = add_textbox(Inches(0.5), y_footer + Inches(0.23), Inches(7.5), Inches(0.18))
tf = txBox.text_frame
set_text(tf, '📅 每 Session 4节课（每周2次×2周） | 每节课50分钟 | 线上小班 4–8人 | 谷雨中文  www.guyuchinese.com', font_size=7, color=RGBColor(0xFF, 0xFF, 0xE0), alignment=PP_ALIGN.CENTER)

# Save
output_path = '/Users/Huan/projects/summercourse/Math/math_course_flyer.pptx'
prs.save(output_path)
print(f'Flyer saved to: {output_path}')
